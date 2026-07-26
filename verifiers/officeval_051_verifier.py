"""
可编辑流程图复刻.pptx 自动评估脚本
依赖: pip install python-pptx Pillow
"""
import sys
import json

SCRIPT_ID = "051"

import math
from decimal import Decimal, InvalidOperation
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.enum.dml import MSO_THEME_COLOR
import zipfile

# ── helpers ──────────────────────────────────────────────────────────────────

def emu2cm(e): return e / 360000
def rgb(r,g,b): return RGBColor(r,g,b)


def normalize_integral_coordinates(prs):
    """规范化整数值小数，并返回不符合整数坐标类型的原始值。"""
    invalid_values = []
    node_attributes = (
        ('.//a:off | .//a:chOff | .//a:pt', ('x', 'y')),
        ('.//a:ext | .//a:chExt | .//p:sldSz', ('cx', 'cy')),
        ('.//a:path', ('w', 'h')),
    )
    for part in prs.part.package.iter_parts():
        root = getattr(part, '_element', None)
        if root is None:
            continue
        for xpath, attr_names in node_attributes:
            for node in root.xpath(xpath):
                for attr_name in attr_names:
                    raw_value = node.get(attr_name)
                    if raw_value is None:
                        continue
                    try:
                        value = Decimal(raw_value)
                    except InvalidOperation:
                        invalid_values.append(raw_value)
                        continue
                    if not value.is_finite() or value != value.to_integral_value():
                        invalid_values.append(raw_value)
                        continue
                    normalized = str(int(value))
                    if normalized != raw_value:
                        node.set(attr_name, normalized)
    return invalid_values


def shape_bounds(s):
    """返回标准化后的 (left_cm, top_cm, w_cm, h_cm)，兼容负宽/负高线段。"""
    l1, t1 = emu2cm(s.left), emu2cm(s.top)
    l2, t2 = emu2cm(s.left + s.width), emu2cm(s.top + s.height)
    return min(l1, l2), min(t1, t2), abs(l2 - l1), abs(t2 - t1)

def shape_right(s):
    l,t,w,h = shape_bounds(s)
    return l + w

def shape_bottom(s):
    l,t,w,h = shape_bounds(s)
    return t + h

def in_range(v, lo, hi): return lo <= v <= hi

def get_fill_color(shape):
    """尝试获取形状填充色 RGBColor，失败返回 None"""
    try:
        fill = shape.fill
        if fill.type is None: return None
        fg = fill.fore_color
        if fg.type is not None:
            return fg.rgb
    except Exception:
        pass
    return None

def get_line_color(shape):
    # ① 优先按 srgbClr 直接取
    try:
        c = shape.line.color.rgb
        if c is not None:
            return c
    except Exception:
        pass
    # ② XML 兜底：解析 <a:ln>/<a:solidFill> 下的 schemeClr / sysClr / srgbClr
    #    并按需应用 lumMod/lumOff/shade/tint（仅对灰阶基色应用，避免把
    #    accent 彩色调暗到分类阈值以下）。
    try:
        sp = shape._element.spPr
        ln = sp.find(qn('a:ln')) if sp is not None else None
        if ln is None:
            return None
        sf = ln.find(qn('a:solidFill'))
        if sf is None or len(sf) == 0:
            return None
        clr_node = sf[0]
        master = _master_of_shape(shape)
        base = _resolve_color_node(clr_node, master)
        if base is None:
            return None
        return _apply_color_mods(base, clr_node)
    except Exception:
        return None

def _master_of_shape(shape):
    """尝试拿到 shape 所在 slide 的 master，用于解析 schemeClr。"""
    try:
        part = shape.part
        slide = getattr(part, 'slide', None)
        if slide is not None:
            return slide.slide_layout.slide_master
    except Exception:
        pass
    return None

def _apply_color_mods(rgb, clr_node):
    """把 clr_node 上的 lumMod/lumOff/shade/tint 应用到 rgb。
    仅对灰阶基色（r≈g≈b，如 bg1=白、tx1=黑）应用亮度调整——
    办公软件里对彩色 accent 加 lumMod 只是微调色深，色相不变，
    强行乘算会把绿/蓝/橙压过颜色分类阈值，反而漏判。"""
    if rgb is None or clr_node is None:
        return rgb
    r, g, b = rgb[0], rgb[1], rgb[2]
    is_gray = max(r, g, b) - min(r, g, b) <= 8
    if not is_gray:
        return rgb
    def _val(tag):
        n = clr_node.find(qn('a:' + tag))
        if n is None: return None
        v = n.get('val')
        try:
            return int(v) / 100000.0 if v is not None else None
        except Exception:
            return None
    lm = _val('lumMod'); lo = _val('lumOff')
    sh = _val('shade');  ti = _val('tint')
    fr, fg, fb = float(r), float(g), float(b)
    if lm is not None:
        fr *= lm; fg *= lm; fb *= lm
    if lo is not None:
        fr += 255 * lo; fg += 255 * lo; fb += 255 * lo
    if sh is not None:
        fr *= (1 - sh); fg *= (1 - sh); fb *= (1 - sh)
    if ti is not None:
        fr = fr + (255 - fr) * ti
        fg = fg + (255 - fg) * ti
        fb = fb + (255 - fb) * ti
    def clamp(x): return max(0, min(255, int(round(x))))
    return RGBColor(clamp(fr), clamp(fg), clamp(fb))

def line_color_is_black_or_dark(shape):
    """判定形状线条颜色为"黑/深色"，兼容三种情况：
      ① 显式 RGB 值为深色（走 color_is_black_or_dark）；
      ② 采用主题色 schemeClr 且槽位为 tx1/dk1/dk2（办公软件默认渲染为黑/深色文本槽）；
      ③ <a:ln> 中 solidFill/gradFill/pattFill 都不存在，或存在空 <a:solidFill/>（无子色元素，
         继承主题默认线条色，PowerPoint 默认为黑）。
    """
    c = get_line_color(shape)
    if c is not None and color_is_black_or_dark(c):
        return True
    try:
        sp = shape._element.spPr
        ln = sp.find(qn('a:ln')) if sp is not None else None
        if ln is None:
            # 无显式 a:ln：继承主题（默认黑）
            return True
        # 空的 <a:solidFill/>（无 srgbClr/schemeClr 子元素）等价于未指定色 → 继承黑
        sf = ln.find(qn('a:solidFill'))
        if sf is not None and len(sf) == 0:
            return True
        # 检查 schemeClr：tx1/dk1/dk2 皆为深色文本槽
        for fill_tag in ('a:solidFill', 'a:gradFill', 'a:pattFill', 'a:noFill'):
            f = ln.find(qn(fill_tag))
            if f is not None:
                if fill_tag == 'a:noFill':
                    return False
                if fill_tag == 'a:solidFill':
                    sc = f.find(qn('a:schemeClr'))
                    if sc is not None and sc.get('val') in ('tx1', 'dk1', 'dk2'):
                        return True
                    # 其余情况（srgbClr 等）交由 RGB 分支处理
                return False
        # <a:ln> 存在但未指定 fill：继承主题默认（黑）
        return True
    except Exception:
        return False

def get_line_width_pt(shape):
    try:
        w = shape.line.width
        if w is None: return None
        return w / 12700  # EMU -> pt
    except Exception:
        return None

def line_width_in_range(shape, lo, hi):
    """判定形状线宽是否落在 [lo, hi] 磅区间。
    若形状未显式指定线宽（继承主题/默认，python-pptx 会返回 0 EMU 或 None），
    按 PowerPoint 默认线宽 0.75 磅代入判定。
    """
    wpt = get_line_width_pt(shape)
    # 显式未设 w 时，python-pptx line.width 通常返回 0 EMU，需回落到默认
    if wpt is None or wpt <= 0:
        wpt = 0.75   # PowerPoint 默认线宽
    return in_range(wpt, lo, hi)

def line_is_single_solid(shape):
    """边线为"单实线"：线型为实线（prstDash 为 solid 或缺省默认实线），
    且复合线型为单线（cmpd 为 sng 或缺省），排除虚线/双线/粗细线等。
    以办公软件实际渲染的 a:ln 定义为准。"""
    try:
        sp = shape._element.spPr
        ln = sp.find(qn('a:ln')) if sp is not None else None
        if ln is None:
            # 无显式线定义：无法确认存在实线边框
            return False
        cmpd = ln.get('cmpd')
        if cmpd is not None and cmpd != 'sng':
            return False  # dbl/thickThin/thinThick/tri 等复合线，非单线
        prst = ln.find(qn('a:prstDash'))
        if prst is not None and prst.get('val') not in (None, 'solid'):
            return False  # 明确的虚线/点线
        if ln.find(qn('a:custDash')) is not None:
            return False  # 自定义虚线
        if ln.find(qn('a:noFill')) is not None:
            return False  # 无线
        return True
    except Exception:
        return False

def is_dash_line(shape):
    """检测线型是否为虚线（支持真正dash_style，也支持短线段拼虚线外框）。
    办公软件中"虚线"涵盖 dash / dot / dashDot 等所有非实线模式，
    python-pptx 的 MSO_LINE_DASH_STYLE 枚举名可能是 DASH/ROUND_DOT/SQUARE_DOT/DASH_DOT/…，
    也可能是 SYSTEM_DASH/SYSTEM_DOT 等，均视为虚线。"""
    try:
        dash = shape.line.dash_style
        if dash is not None:
            name = str(dash).upper()
            if any(k in name for k in ('DASH', 'DOT')):
                return True
    except Exception:
        pass
    # 直接看 XML：<a:prstDash val="..."/> 非 solid 皆为虚线
    try:
        sp = shape._element.spPr
        ln = sp.find(qn('a:ln')) if sp is not None else None
        if ln is not None:
            pd = ln.find(qn('a:prstDash'))
            if pd is not None and (pd.get('val') or 'solid') != 'solid':
                return True
    except Exception:
        pass
    return False

def is_drawn_as_dashes(shape, shapes, max_gap=0.42, min_segments=10):
    """
    检测一个形状是否被大量相邻短线段"拼"出了虚线外框。
    必须在该形状自身四条边附近都有同色短线段，避免把附近外层虚线框
    误判成内部实线矩形的虚线边框。
    """
    l,t,w,h = shape_bounds(shape)
    c = get_line_color(shape)
    if c is None: return False
    tol = 0.12
    same = [s for s in shapes if s is not shape and get_line_color(s) == c]
    short_segs = []
    for s in same:
        sl,st,sw,sh = shape_bounds(s)
        if max(sw, sh) <= max_gap and min(sw, sh) <= 0.03:
            short_segs.append((sl, st, sw, sh))

    top = [x for x in short_segs if x[2] > x[3] and abs(x[1]-t) <= tol and l-tol <= x[0] <= l+w+tol]
    bottom = [x for x in short_segs if x[2] > x[3] and abs(x[1]-(t+h)) <= tol and l-tol <= x[0] <= l+w+tol]
    left = [x for x in short_segs if x[3] > x[2] and abs(x[0]-l) <= tol and t-tol <= x[1] <= t+h+tol]
    right = [x for x in short_segs if x[3] > x[2] and abs(x[0]-(l+w)) <= tol and t-tol <= x[1] <= t+h+tol]
    return (len(top) + len(bottom) + len(left) + len(right) >= min_segments
            and all(len(side) >= 2 for side in (top, bottom, left, right)))

def fill_is_none_or_transparent(shape):
    """检查形状是否无填充/透明填充；浅色填充不算无填充。"""
    try:
        sp_pr = shape._element.spPr
        if sp_pr.find(qn('a:noFill')) is not None:
            return True
        solid_fill = sp_pr.find(qn('a:solidFill'))
        if solid_fill is not None:
            alpha = solid_fill.find('.//' + qn('a:alpha'))
            return alpha is not None and alpha.get('val') == '0'
        fill = shape.fill
        return fill.type is None
    except Exception:
        return False

def fill_is_visually_empty(shape):
    """办公软件中"内部无填充"的有效判定：真正无填充 / 完全透明，
    或填充为近白色（在白底页面上渲染后与页面无法区分，视觉上等同无填充）。"""
    if fill_is_none_or_transparent(shape):
        return True
    c = get_fill_color(shape)
    if c is None:
        return True
    return color_is_white_or_near(c)

def color_is_light_gray(c):
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return abs(r-g) <= 35 and abs(g-b) <= 35 and 120 <= r <= 230 and 120 <= g <= 230 and 120 <= b <= 230

def color_is_light(c):
    if c is None: return True
    r,g,b = c[0],c[1],c[2]
    lum = 0.299*r + 0.587*g + 0.114*b
    return lum > 200

def color_is_white_or_grayish_white(c):
    """白色或极浅灰白色：高亮度且接近中性（低彩度）。None 视为办公软件默认白底。"""
    if c is None:
        return True
    r, g, b = c[0], c[1], c[2]
    return min(r, g, b) >= 224 and (max(r, g, b) - min(r, g, b)) <= 12

def color_is_white_or_near(c):
    if c is None: return True
    r,g,b = c[0],c[1],c[2]
    return r>220 and g>220 and b>220

def color_is_cyan_green(c):
    """青绿色/深青绿色 ~ teal"""
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return (g >= 140 and b >= 120 and r < 170) or (r < 80 and g >= 80 and b >= 80 and abs(g-b) <= 60)

def color_is_light_cyan(c):
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return r>180 and g>220 and b>220

def color_is_purple(c):
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return r>100 and b>100 and g < min(r,b)-20

def color_is_orange(c):
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return r>140 and g>50 and b<140 and r>g+25

def color_is_blue(c):
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return b>150 and b > r+30 and b > g+30

def color_is_green(c):
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return g>150 and g > r+30 and g > b+20

def color_is_pink(c):
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return r>180 and b>120 and g<180 and r>b+20

def color_is_black_or_dark(c):
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return r<80 and g<80 and b<80

def color_is_light_orange(c):
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return r>220 and g>180 and b>150 and r>b+50

def color_is_orangey_light_fill(c):
    """办公软件中实际呈现为"浅橙色"的填充：从明显浅橙到暖调奶油近白都算。
    判据：暖色相（r>=g>=b 且带暖偏移 r>b），亮度较高。以肉眼所见为准。"""
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    if color_is_light_orange(c):
        return True
    # 暖调浅色（含奶油/浅橙近白）：红≥绿≥蓝且明显偏暖、整体偏亮
    return r >= g >= b and (r - b) >= 4 and r >= 235 and g >= 210 and b >= 190

def color_is_light_blue(c):
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return b>180 and r>180 and g>180 and b>=r

def color_is_light_green(c):
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return g>200 and r>180 and b>180

def color_is_light_purple(c):
    if c is None: return False
    r,g,b = c[0],c[1],c[2]
    return r>180 and b>180 and g>180 and (r+b) > g*2+40

def font_is_arial_calibri(run_font):
    if run_font is None: return False
    name = (run_font or '').lower()
    return 'arial' in name or 'calibri' in name

def has_smaller_suffix_text(shape, prefix, suffix):
    """检测类似 Lrec/Lreg：suffix（如 rec）在办公软件中显示得比 prefix（L）小。
    覆盖所有可能的实现方式：
      ① suffix 与 prefix 分属不同 run，且 suffix run 字号更小；
      ② suffix 设为下标/上标（baseline 属性），办公软件会自动缩小显示；
      ③ 同一 run 内但 suffix 字符字号更小（罕见，仍按 run 字号比较）。"""
    from pptx.oxml.ns import qn as _qn
    runs = [r for r in get_text_runs(shape) if r.text.strip()]
    if not runs:
        return False
    text = ''.join(r.text for r in runs).replace('_', '').replace(' ', '')
    target = prefix + suffix
    if target.lower() not in text.lower():
        return False

    prefix_sizes = []
    suffix_sizes = []
    suffix_has_baseline = False
    for r in runs:
        sz = r.font.size.pt if r.font.size else None
        # 读取 baseline（下标/上标）——有则视觉缩小
        rpr = r._r.find(_qn('a:rPr'))
        baseline = rpr.get('baseline') if rpr is not None else None
        rtext = r.text.replace('_', '')
        for ch in rtext:
            if ch.lower() == prefix.lower():
                if sz is not None:
                    prefix_sizes.append(sz)
            elif ch.lower() in suffix.lower():
                if sz is not None:
                    suffix_sizes.append(sz)
                if baseline not in (None, '', '0'):   # ② 下标/上标
                    suffix_has_baseline = True
    # ② 后缀被设为下标/上标 → 办公软件显示更小
    if suffix_has_baseline and (prefix.lower() in ''.join(r.text.lower() for r in runs)):
        return True
    # ①③ 按字号比较：后缀最大字号 < 前缀最小字号
    return bool(prefix_sizes and suffix_sizes and max(suffix_sizes) < min(prefix_sizes))

def color_is_genuine_light_blue(c):
    """真正带浅蓝色调的填充（区别于纯白）：整体偏亮，且蓝分量明显高于红/绿。
    覆盖办公软件里"浅蓝色"的各种呈现，但排除纯白/近白（无蓝调）。"""
    if c is None:
        return False
    r, g, b = c[0], c[1], c[2]
    if not (b >= 200 and b >= r and b >= g):
        return False
    # 需有可见蓝调：蓝比红或绿至少高一点，避免把纯白(相等)判成浅蓝
    return (b - min(r, g)) >= 6


def get_text_runs(shape):
    runs = []
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                runs.append(run)
    except Exception:
        pass
    return runs

def shape_text(shape):
    try:
        return shape.text_frame.text
    except Exception:
        return ''

def shapes_on_slide(slide):
    result = []
    def collect(sp_list):
        for s in sp_list:
            result.append(s)
            try:
                collect(s.shapes)
            except Exception:
                pass
    collect(slide.shapes)
    return result

def is_rounded_rect(shape):
    try:
        sp_pr = shape._element.spPr
        prstGeom = sp_pr.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prst = prstGeom.get('prst', '') or ''
            return 'roundRect' in prst or 'round' in prst.lower()
    except Exception:
        pass
    return False

def is_line_or_connector(shape):
    try:
        return shape.shape_type in (6, 9)  # LINE, CONNECTOR (freeform connector)
    except Exception:
        return False

def get_connector_dash(shape):
    try:
        ln = shape.line
        return ln.dash_style
    except Exception:
        return None


def center(shape):
    l,t,w,h = shape_bounds(shape)
    return l+w/2, t+h/2

def freeform_points(shape):
    """读取 FREEFORM 形状 <a:custGeom>/<a:pathLst>/<a:path> 下的顶点坐标，
    按形状实际位置/尺寸换算成 cm 坐标列表 [(x_cm, y_cm), ...]。
    custGeom 内的坐标以 path 自身声明的宽高（w/h，单位EMU的“路径坐标系”）为参照，
    需按 shape 实际宽高做比例缩放后再加上 shape 左上角偏移。失败返回 None。"""
    try:
        sp_pr = shape._element.spPr
        cg = sp_pr.find(qn('a:custGeom'))
        if cg is None:
            return None
        path_lst = cg.find(qn('a:pathLst'))
        if path_lst is None:
            return None
        path = path_lst.find(qn('a:path'))
        if path is None:
            return None
        pw = int(path.get('w', 0)) or 1
        ph = int(path.get('h', 0)) or 1
        l, t, w, h = shape_bounds(shape)
        pts = []
        for child in path:
            tag = etree_local_tag(child)
            if tag in ('moveTo', 'lnTo'):
                pt = child.find(qn('a:pt'))
                if pt is None:
                    continue
                px, py = int(pt.get('x')), int(pt.get('y'))
                pts.append((l + (px / pw) * w, t + (py / ph) * h))
            elif tag == 'cubicBezTo':
                for pt in child.findall(qn('a:pt')):
                    px, py = int(pt.get('x')), int(pt.get('y'))
                    pts.append((l + (px / pw) * w, t + (py / ph) * h))
        return pts if pts else None
    except Exception:
        return None

def etree_local_tag(el):
    tag = el.tag
    return tag.split('}', 1)[1] if '}' in tag else tag

def bbox_intersection(a, b):
    ax,ay,aw,ah = shape_bounds(a); bx,by,bw,bh = shape_bounds(b)
    x1,y1 = max(ax,bx), max(ay,by)
    x2,y2 = min(ax+aw,bx+bw), min(ay+ah,by+bh)
    if x2 <= x1 or y2 <= y1: return 0.0
    return (x2-x1)*(y2-y1)

def bbox_area(shape):
    _,_,w,h = shape_bounds(shape)
    return max(0,w)*max(0,h)

def overlaps(a,b,ratio=0.35):
    inter = bbox_intersection(a,b)
    if inter <= 0: return False
    return inter / max(0.01, min(bbox_area(a), bbox_area(b))) >= ratio

def box_in(shape, left=None, top=None, width=None, height=None, right=None, bottom=None):
    l,t,w,h = shape_bounds(shape)
    checks = []
    if left: checks.append(in_range(l, *left))
    if top: checks.append(in_range(t, *top))
    if width: checks.append(in_range(w, *width))
    if height: checks.append(in_range(h, *height))
    if right: checks.append(in_range(l+w, *right))
    if bottom: checks.append(in_range(t+h, *bottom))
    return all(checks)

def canonical_text(s):
    return ''.join(ch.lower() for ch in s.replace('_', '').replace('／', '/') if not ch.isspace())

def text_contains(shape, needle):
    raw = shape_text(shape)
    compact_raw = canonical_text(raw)
    compact_need = canonical_text(needle)
    alnum_raw = ''.join(ch for ch in compact_raw if ch.isalnum())
    alnum_need = ''.join(ch for ch in compact_need if ch.isalnum())
    return needle.lower() in raw.lower() or compact_need in compact_raw or (alnum_need and alnum_need in alnum_raw)

def text_norm(shape):
    return ' '.join(shape_text(shape).replace('\r','\n').split())

def line_count(shape):
    txt = shape_text(shape).replace('\r', '\n')
    return len([x for x in txt.split('\n') if x.strip()])


# ── 文本行数：按办公软件实际渲染估算 ──────────────────────────────────
# 细则里的"分几行"是以办公软件打开时实际看到的行数为准，而不是源 XML 里
# 段落数。python-pptx 的 line_count 只数 \n/段落，会把"文本超出文本框宽度
# 被自动换行"的情形误判为 1 行。这里改用字体度量+文本框可用宽度模拟换行。

_FONT_PATHS = {
    ('arial', True): 'C:/Windows/Fonts/arialbd.ttf',
    ('arial', False): 'C:/Windows/Fonts/arial.ttf',
    ('calibri', True): 'C:/Windows/Fonts/calibrib.ttf',
    ('calibri', False): 'C:/Windows/Fonts/calibri.ttf',
    ('times new roman', True): 'C:/Windows/Fonts/timesbd.ttf',
    ('times new roman', False): 'C:/Windows/Fonts/times.ttf',
}
_FONT_CACHE = {}

def _load_font(name, bold, sz_pt):
    key = ((name or 'arial').lower().strip(), bool(bold), round(float(sz_pt), 2))
    if key in _FONT_CACHE:
        return _FONT_CACHE[key]
    try:
        from PIL import ImageFont
    except Exception:
        _FONT_CACHE[key] = None
        return None
    path = _FONT_PATHS.get(key[:2])
    if path is None:
        # 未知字体退回 Arial（保持度量可用；实际字宽差异不大）
        path = _FONT_PATHS[('arial', bool(bold))]
    sz_px = max(1, int(round(sz_pt * 96 / 72)))
    try:
        f = ImageFont.truetype(path, sz_px)
    except Exception:
        try:
            f = ImageFont.load_default()
        except Exception:
            f = None
    _FONT_CACHE[key] = f
    return f

def _measure_text_cm(text, font_name, bold, sz_pt):
    f = _load_font(font_name, bold, sz_pt)
    if f is None or not text:
        return 0.0
    try:
        bbox = f.getbbox(text)
    except Exception:
        return 0.0
    return (bbox[2] - bbox[0]) / 96.0 * 2.54

def _runs_of_paragraph(p):
    runs = []
    for r in p.runs:
        name = r.font.name or 'Arial'
        try:
            sz = r.font.size.pt if r.font.size else None
        except Exception:
            sz = None
        if sz is None:
            sz = 18.0  # 未显式声明字号时按 PPT 默认体量估算
        bold = bool(r.font.bold)
        runs.append((r.text, name, bold, sz))
    return runs

def _paragraph_wrap_lines(runs, avail_cm, tolerance=1.05):
    """段落内按空格贪心换行估算行数。tolerance 给 PPT 更紧的字距一点余地。
    额外处理：单个词宽度超过可用宽度时，PowerPoint 会字符级断行，此时用
    ceil(段总宽 / 可用宽度) 作为该段最少行数。"""
    import math
    if not runs:
        return 1
    # 段总宽 & 是否存在超长词
    total_w = 0.0
    max_token_w = 0.0
    for text, name, bold, sz in runs:
        if not text:
            continue
        total_w += _measure_text_cm(text, name, bold, sz)
        for part in text.split(' '):
            if part:
                w = _measure_text_cm(part, name, bold, sz)
                if w > max_token_w:
                    max_token_w = w
    avail_effective = avail_cm * tolerance
    # 拼接文本、按空格分词；每个词度量宽度用其所属 run 的字体
    tokens = []  # (text, name, bold, sz)
    for text, name, bold, sz in runs:
        if not text:
            continue
        parts = text.split(' ')
        for i, part in enumerate(parts):
            if part:
                tokens.append((part, name, bold, sz))
            if i < len(parts) - 1:
                tokens.append((' ', name, bold, sz))
    if not tokens:
        return 1
    lines = 1
    cur_w = 0.0
    for text, name, bold, sz in tokens:
        w = _measure_text_cm(text, name, bold, sz)
        if cur_w == 0.0:
            if text == ' ':
                continue
            cur_w = w
            continue
        if cur_w + w <= avail_effective:
            cur_w += w
        else:
            lines += 1
            cur_w = 0.0 if text == ' ' else w
    # 单词强制字符级断行下限（避免"某个词本身超宽却只算 1 行"的漏判）
    if max_token_w > avail_effective:
        lines = max(lines, int(math.ceil(total_w / avail_effective)))
    return lines

def rendered_line_count(shape):
    """按办公软件实际渲染估算文本行数。"""
    try:
        body = shape._element.find('.//' + qn('a:bodyPr'))
        # wrap="none" 时不换行 —— 段落数即行数
        wrap = 'square'
        l_ins = r_ins = 91440  # PPT 默认 lIns/rIns 单位为 EMU；缺省 0.1 inch
        font_scale = 1.0
        if body is not None:
            wrap = body.get('wrap', 'square') or 'square'
            l_ins = int(body.get('lIns', str(l_ins)) or l_ins)
            r_ins = int(body.get('rIns', str(r_ins)) or r_ins)
            # "根据形状自动调整文字大小"：<a:normAutofit fontScale="90000"/> 表示办公软件
            # 实际渲染时把字号整体缩放到 90%（单位千分之一），文字宽度须按此折算，
            # 否则会把已被自动缩小、实际能容纳的一行误判为换行后的多行。
            autofit = body.find(qn('a:normAutofit'))
            if autofit is not None:
                fs = autofit.get('fontScale')
                if fs:
                    try:
                        font_scale = int(fs) / 100000.0
                    except Exception:
                        font_scale = 1.0
        _, _, w_cm, _ = shape_bounds(shape)
        inset_cm = (l_ins + r_ins) / 914400.0 * 2.54
        avail_cm = max(0.05, w_cm - inset_cm)
        paragraphs = list(shape.text_frame.paragraphs)
        if not paragraphs:
            return 0
        total = 0
        any_visible = False
        for p in paragraphs:
            runs = _runs_of_paragraph(p)
            if font_scale != 1.0:
                runs = [(t, n, b, sz * font_scale) for (t, n, b, sz) in runs]
            text = ''.join(r[0] for r in runs)
            if not text.strip():
                # 空段落 —— PPT 里也会显示一行空行，但细则中"分几行"通常忽略
                continue
            any_visible = True
            if wrap == 'none':
                total += 1
                continue
            # 段落内全文宽度
            full_w = sum(_measure_text_cm(t, n, b, s) for (t, n, b, s) in runs)
            if full_w <= avail_cm * 1.05:
                total += 1
            else:
                total += _paragraph_wrap_lines(runs, avail_cm)
        return total if any_visible else 0
    except Exception:
        # 度量失败退回原始段落数
        return line_count(shape)

def text_font_ok(shape, fonts=('arial','calibri'), size=None, bold=None, color_pred=None, center_align=None):
    """按 run 检查文本属性。默认严格模式：
      - bold=True 时，run 必须显式声明 b="1"（bold is None 视作未加粗）；
      - bold=False 时，任何 run 显式 b="1" 都视作失败；
      - 字体/字号/颜色仍允许缺失继承 —— 因为主题里通常会强制指定，
        且这些属性在办公软件里可用系统默认渲染而不影响细则判定。
    fonts=None 时不校验字体（用于细则未限定字体的场景）。"""
    runs = [r for r in get_text_runs(shape) if r.text.strip()]
    if not runs:
        return False
    font_hits = 0; size_hits = 0; color_hits = 0
    checked_font = checked_size = checked_color = 0
    bold_explicit_true = 0
    bold_explicit_false = 0
    for r in runs:
        name = (r.font.name or '').lower()
        if fonts is not None and name:
            checked_font += 1
            if any(f in name for f in fonts): font_hits += 1
        if r.font.size:
            checked_size += 1
            pt = r.font.size.pt
            if size is None or in_range(pt, *size): size_hits += 1
        # 加粗：严格计数——仅 XML 显式声明才算数
        if r.font.bold is True:
            bold_explicit_true += 1
        elif r.font.bold is False:
            bold_explicit_false += 1
        try:
            c = r.font.color.rgb
            if c is not None:
                checked_color += 1
                if color_pred is None or color_pred(c): color_hits += 1
        except Exception:
            pass
    if checked_font and font_hits == 0: return False
    if size is not None and checked_size and size_hits == 0: return False
    if bold is True:
        # 严格加粗：所有可见 run 必须显式 b="1"
        if bold_explicit_true == 0 or bold_explicit_true < len(runs):
            return False
    elif bold is False:
        if bold_explicit_true > 0:
            return False
    if color_pred is not None and checked_color and color_hits == 0: return False
    if center_align is not None:
        try:
            aligns = [p.alignment for p in shape.text_frame.paragraphs if p.text.strip()]
            # 2 is CENTER in python-pptx enum; some inherited alignment may be None.
            if any(a is not None for a in aligns) and not any(str(a).upper().endswith('CENTER') or a == 2 for a in aligns):
                return False
        except Exception:
            pass
    return True

def has_arrow_head(shape):
    try:
        ln = shape._element.spPr.find(qn('a:ln'))
        if ln is None: return False
        for tag in ('a:tailEnd', 'a:headEnd'):
            el = ln.find(qn(tag))
            if el is not None and el.get('type', 'triangle') != 'none':
                return True
    except Exception:
        return False

def vertical_arrow_points_down(shape):
    """竖向箭头是否朝下（按办公软件实际几何：结合两端箭头 + flipV 翻转）。
    line 的起点=上端(headEnd)、终点=下端(tailEnd)；flipV 会上下互换。"""
    try:
        ln = shape._element.spPr.find(qn('a:ln'))
        if ln is None:
            return False
        head = ln.find(qn('a:headEnd'))
        tail = ln.find(qn('a:tailEnd'))
        head_arrow = head is not None and head.get('type', 'triangle') != 'none'
        tail_arrow = tail is not None and tail.get('type', 'triangle') != 'none'
        if not (head_arrow or tail_arrow):
            return False
        xfrm = shape._element.find('.//' + qn('a:xfrm'))
        flip = xfrm is not None and xfrm.get('flipV') == '1'
        # 未翻转：tail 在下端；翻转：head 变到下端
        return (tail_arrow and not flip) or (head_arrow and flip)
    except Exception:
        return False

def horizontal_arrow_points_right(shape):
    """水平箭头是否朝右（按办公软件实际几何：结合线段两端箭头 + flipH 翻转）。
    line 的起点=左端(headEnd)、终点=右端(tailEnd)；flipH 会左右互换。"""
    try:
        ln = shape._element.spPr.find(qn('a:ln'))
        if ln is None:
            return False
        head = ln.find(qn('a:headEnd'))
        tail = ln.find(qn('a:tailEnd'))
        head_arrow = head is not None and head.get('type', 'triangle') != 'none'
        tail_arrow = tail is not None and tail.get('type', 'triangle') != 'none'
        if not (head_arrow or tail_arrow):
            return False
        xfrm = shape._element.find('.//' + qn('a:xfrm'))
        flip = xfrm is not None and xfrm.get('flipH') == '1'
        # 未翻转：tail 在右端；翻转：head 变到右端
        right_end_arrow = (tail_arrow and not flip) or (head_arrow and flip)
        return right_end_arrow
    except Exception:
        return False

def is_connector(shape):
    try:
        return shape.shape_type in (MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM)
    except Exception:
        return is_line_or_connector(shape)

def is_horizontal(shape):
    l,t,w,h = shape_bounds(shape)
    return w >= h*2

def is_vertical(shape):
    l,t,w,h = shape_bounds(shape)
    return h >= w*2

def shape_name_or_prst(shape):
    try:
        prstGeom = shape._element.spPr.find(qn('a:prstGeom'))
        prst = prstGeom.get('prst', '') if prstGeom is not None else ''
    except Exception:
        prst = ''
    return (getattr(shape, 'name', '') + ' ' + prst).lower()

def looks_like_rect(shape):
    n = shape_name_or_prst(shape)
    return 'rect' in n or 'rectangle' in n or is_rounded_rect(shape)

def count_shapes_in(box_shape, shapes, pred=lambda s: True):
    l,t,w,h = shape_bounds(box_shape)
    r,b = l+w,t+h
    cnt = 0
    inside = []
    for s in shapes:
        if s is box_shape: continue
        cx,cy = center(s)
        if l <= cx <= r and t <= cy <= b and pred(s):
            cnt += 1; inside.append(s)
    return cnt, inside

def _theme_of_master(master):
    """取母版关联的主题 XML 根节点。"""
    try:
        for rel in master.part.rels.values():
            if 'theme' in rel.reltype:
                return parse_xml(rel.target_part.blob)
    except Exception:
        pass
    return None

def _scheme_color_rgb(name, master):
    """按 clrMap 把占位色名映射到主题 clrScheme，返回 RGBColor 或 None。"""
    if master is None or name is None:
        return None
    try:
        cm = master._element.find(qn('p:clrMap'))
        mapping = dict(cm.attrib) if cm is not None else {}
        mapped = mapping.get(name, name)
        theme = _theme_of_master(master)
        if theme is None:
            return None
        cs = theme.find(qn('a:themeElements') + '/' + qn('a:clrScheme'))
        if cs is None:
            return None
        node = cs.find(qn('a:' + mapped))
        if node is None or len(node) == 0:
            return None
        ch = node[0]
        if ch.tag == qn('a:srgbClr'):
            return RGBColor.from_string(ch.get('val'))
        if ch.tag == qn('a:sysClr'):
            last = ch.get('lastClr')
            return RGBColor.from_string(last) if last else None
    except Exception:
        return None
    return None

def _resolve_color_node(clr_node, master):
    """把 a:srgbClr / a:sysClr / a:schemeClr 解析为 RGBColor。"""
    if clr_node is None:
        return None
    try:
        tag = clr_node.tag
        if tag == qn('a:srgbClr'):
            return RGBColor.from_string(clr_node.get('val'))
        if tag == qn('a:sysClr'):
            last = clr_node.get('lastClr')
            return RGBColor.from_string(last) if last else None
        if tag == qn('a:schemeClr'):
            return _scheme_color_rgb(clr_node.get('val'), master)
    except Exception:
        return None
    return None

def _bg_element(part):
    """取某个部件（幻灯片/版式/母版）cSld 下的 p:bg 节点。"""
    try:
        csld = part._element.find(qn('p:cSld'))
        if csld is None:
            return None
        return csld.find(qn('p:bg'))
    except Exception:
        return None

def _resolve_bgref(bgRef, master):
    """解析 p:bgRef（主题 fillStyle 引用）为 (kind, rgb)。"""
    try:
        idx = int(bgRef.get('idx', '0'))
    except Exception:
        idx = 0
    phclr = _resolve_color_node(bgRef[0], master) if len(bgRef) else None
    theme = _theme_of_master(master)
    if theme is None:
        return ('ref', phclr)
    fmt = theme.find(qn('a:themeElements') + '/' + qn('a:fmtScheme'))
    if fmt is None:
        return ('ref', phclr)
    if idx >= 1001:
        lst = fmt.find(qn('a:bgFillStyleLst')); i = idx - 1000
    else:
        lst = fmt.find(qn('a:fillStyleLst')); i = idx
    if lst is None or i < 1 or i > len(lst):
        return ('ref', phclr)
    fill = lst[i - 1]
    if fill.tag == qn('a:solidFill'):
        node = fill[0] if len(fill) else None
        if node is not None and node.tag == qn('a:schemeClr') and node.get('val') == 'phClr':
            return ('solid', phclr)
        return ('solid', _resolve_color_node(node, master))
    if fill.tag == qn('a:gradFill'):
        return ('gradient', None)
    if fill.tag == qn('a:blipFill'):
        return ('image', None)
    return ('ref', phclr)

def shape_fill_is_semitransparent(shape):
    """填充/图片存在 5%~95% 之间的透明度（明显半透明覆盖），完全透明或不透明都不算。"""
    try:
        sp = shape._element.spPr
        if sp is None:
            return False
        for al in sp.findall('.//' + qn('a:alpha')):
            v = al.get('val')
            if v is None:
                continue
            if 5000 <= int(v) <= 95000:
                return True
    except Exception:
        pass
    return False

def first_run_rgb(shape):
    """取形状内首个有显式颜色的文本 run 的 RGB。"""
    for r in get_text_runs(shape):
        try:
            if r.font.color is not None and r.font.color.type is not None:
                return r.font.color.rgb
        except Exception:
            pass
    return None


class Evaluator:
    def __init__(self, ppt_path):
        self.path = Path(ppt_path)
        self.presentation = None
        self.slide = None
        self.shapes = []
        self.slide_w = 0
        self.slide_h = 0
        self.total: int = 0
        self.hits: list[tuple[int, str, str]] = []
        self.misses: list[tuple[int, str, str]] = []
        self.key = {}
        self.invalid_coordinate_values: list[str] = []

    def load(self):
        self.presentation = Presentation(str(self.path))
        self.invalid_coordinate_values = normalize_integral_coordinates(
            self.presentation
        )
        self.slide_w = emu2cm(self.presentation.slide_width)
        self.slide_h = emu2cm(self.presentation.slide_height)
        if len(self.presentation.slides):
            self.slide = self.presentation.slides[0]
            self.shapes = shapes_on_slide(self.slide)

    def add(self, points: int, desc: str, ok: object, detail: str = ''):
        if bool(ok):
            self.total += points
            self.hits.append((points, desc, detail))
        else:
            self.misses.append((points, desc, detail))

    def find_shapes(self, pred):
        return [s for s in self.shapes if pred(s)]

    def find_text(self, needle, area=None):
        candidates = [s for s in self.shapes if text_contains(s, needle)]
        if area:
            candidates = [s for s in candidates if box_in(s, **area)]
        return candidates

    def dimension1(self):
        reasons = []
        if self.path.suffix.lower() != '.pptx':
            reasons.append('交付文件不是 .pptx')
            return False, reasons
        try:
            self.load()
        except Exception as e:
            reasons.append(f'文件无法正常打开: {e}')
            return False, reasons
        if self.invalid_coordinate_values:
            samples = ', '.join(self.invalid_coordinate_values[:3])
            reasons.append(
                f'PPT XML 包含非整数坐标（共 {len(self.invalid_coordinate_values)} 处，'
                f'示例：{samples}）'
            )
        if len(self.presentation.slides) != 1:
            reasons.append(f'交付PPT不是1页，而是 {len(self.presentation.slides)} 页')
        if not self.slide:
            reasons.append('没有可检查的幻灯片')
            return False, reasons

        return len(reasons) == 0, reasons

    def detect_key_objects(self):
        sh = self.shapes
        def rrects(lo=None, to=None, wo=None, ho=None):
            return [s for s in sh if looks_like_rect(s) and is_rounded_rect(s)
                    and (lo is None or in_range(emu2cm(s.left),*lo))
                    and (to is None or in_range(emu2cm(s.top),*to))
                    and (wo is None or in_range(abs(emu2cm(s.width)),*wo))
                    and (ho is None or in_range(abs(emu2cm(s.height)),*ho))]
        self.key['left_box']   = (rrects(lo=(0.5,5.4),  to=(4.3,11.8),  wo=(4.2,5.5),   ho=(5.5,7.2)) or [None])[0]
        self.key['seq_box']    = (rrects(lo=(7.8,14.4),  to=(1.2,7.0),   wo=(5.5,6.7),   ho=(3.6,5.6)) or [None])[0]
        self.key['purple_box'] = (rrects(lo=(7.2,13.8),  to=(7.5,12.8),  wo=(5.0,6.0),   ho=(4.7,5.6)) or [None])[0]
        self.key['orange_box'] = (rrects(lo=(15.1,29.7), to=(1.2,9.3),   wo=(13.0,15.0), ho=(6.0,8.0)) or [None])[0]
        self.key['green_box']  = (rrects(lo=(29.2,36.5), to=(1.6,11.0),  wo=(4.5,5.5),   ho=(6.8,8.3)) or [None])[0]
        self.key['pink_box']   = (rrects(lo=(29.5,36.5), to=(10.5,16.5), wo=(4.5,5.3),   ho=(4.5,5.5)) or [None])[0]
        self.key['blue_box']   = (rrects(lo=(15.0,25.3), to=(10.0,16.0), wo=(8.0,9.5),   ho=(4.4,5.0)) or [None])[0]
        self.key['legend_box'] = (rrects(lo=(8.8,29.0),  to=(16.5,18.7), wo=(18.5,19.5), ho=(1.3,1.6)) or [None])[0]

    def line_candidates(self, area=None, color_pred=None, width=None, dashed=None, horizontal=None, vertical=None, arrow=None):
        result = []
        for s in self.shapes:
            if not (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s)):
                continue
            if area:
                l,t,w,h = shape_bounds(s)
                al = area.get('left', (-999,999)); at = area.get('top', (-999,999))
                aw = area.get('width', (-999,999)); ah = area.get('height', (-999,999))
                # 线段常有0宽/0高，因此区域主要按左上坐标/长度容差判断
                if not (in_range(l, *al) and in_range(t, *at) and in_range(w, *aw) and in_range(h, *ah)):
                    continue
            c = get_line_color(s)
            if color_pred and not color_pred(c): continue
            wpt = get_line_width_pt(s)
            if width and wpt is not None and not in_range(wpt, *width): continue
            if dashed is not None and is_dash_line(s) != dashed: continue
            if horizontal is not None and is_horizontal(s) != horizontal: continue
            if vertical is not None and is_vertical(s) != vertical: continue
            if arrow is not None and has_arrow_head(s) != arrow: continue
            result.append(s)
        return result

    def rounded_rects(self, area=None, line_color=None, fill_color=None, width=None, dashed=None):
        result = []
        for s in self.shapes:
            if not (looks_like_rect(s) and is_rounded_rect(s)): continue
            if area and not box_in(s, **area): continue
            if line_color and not line_color(get_line_color(s)): continue
            if fill_color and not fill_color(get_fill_color(s)): continue
            lw = get_line_width_pt(s)
            if width and lw is not None and not in_range(lw, *width): continue
            if dashed is not None and self.is_dashed_outline(s) != dashed: continue
            result.append(s)
        return result

    def is_dashed_outline(self, shape):
        """"虚线矩形"判定：仅承认矩形自身边线的 prstDash（真正的虚线属性）。
        不接受"边线透明实线 + 周围用短线段拼出虚线视觉"这类拼装做法 ——
        细则要求矩形本身是虚线矩形。"""
        return is_dash_line(shape)

    def _master(self):
        try:
            return self.slide.slide_layout.slide_master
        except Exception:
            return None

    def background_kind_and_color(self):
        """按 幻灯片→版式→母版 继承顺序解析当前生效的整页背景。
        返回 (kind, rgb)，kind ∈ {'solid','gradient','image','none'}。"""
        master = self._master()
        try:
            parts = [self.slide, self.slide.slide_layout, master]
        except Exception:
            parts = [self.slide]
        for part in parts:
            if part is None:
                continue
            bg = _bg_element(part)
            if bg is None:
                continue
            bgpr = bg.find(qn('p:bgPr'))
            if bgpr is not None:
                if bgpr.find(qn('a:noFill')) is not None:
                    return ('none', None)
                solid = bgpr.find(qn('a:solidFill'))
                if solid is not None:
                    node = solid[0] if len(solid) else None
                    return ('solid', _resolve_color_node(node, master))
                if bgpr.find(qn('a:gradFill')) is not None:
                    return ('gradient', None)
                if bgpr.find(qn('a:blipFill')) is not None:
                    return ('image', None)
                continue
            bgref = bg.find(qn('p:bgRef'))
            if bgref is not None:
                return _resolve_bgref(bgref, master)
        # 无任何显式背景 → 办公软件默认白色
        return ('solid', RGBColor(0xFF, 0xFF, 0xFF))

    def has_watermark(self):
        """检测页面是否存在水印。办公软件中水印通常表现为：
        1) 背景/大幅图片以半透明形式铺在页面上；
        2) 大面积、半透明或极浅的重复文字（如"样张""机密"）压在内容之上。
        仅识别真正的水印特征，普通浅色装饰不计入。"""
        page_area = self.slide_w * self.slide_h
        for s in self.shapes:
            area = bbox_area(s)
            if area < page_area * 0.25:
                continue
            # 半透明的大幅图片/形状铺满页面 → 图片型水印
            is_pic = getattr(s, 'shape_type', None) == MSO_SHAPE_TYPE.PICTURE
            if (is_pic or area > page_area * 0.5) and shape_fill_is_semitransparent(s):
                return True
            # 大面积文字水印：字号大、半透明或极浅灰色
            txt = shape_text(s).strip()
            if txt:
                sizes = [r.font.size.pt for r in get_text_runs(s) if r.font.size]
                big_font = bool(sizes) and max(sizes) >= 40
                col = first_run_rgb(s)
                faint = shape_fill_is_semitransparent(s) or (col is not None and color_is_light_gray(col))
                if area > page_area * 0.3 and big_font and faint:
                    return True
        return False



    def rect_covering_text(self, text_shape, rect_pred=lambda r: True, margin=0.25):
        """返回包含文本中心点、且满足条件的最小圆角矩形，避免先匹配到外层大框。"""
        cx, cy = center(text_shape)
        matches = []
        for r in self.shapes:
            if r is text_shape or not (looks_like_rect(r) and is_rounded_rect(r)):
                continue
            l,t,w,h = shape_bounds(r)
            if l-margin <= cx <= l+w+margin and t-margin <= cy <= t+h+margin and rect_pred(r):
                matches.append(r)
        if not matches:
            return None
        return min(matches, key=bbox_area)

    def text_in_area(self, needle, area, font_kwargs=None, require_lines=None):
        for s in self.find_text(needle):
            if not box_in(s, **area):
                continue
            if require_lines is not None and rendered_line_count(s) != require_lines:
                continue
            if font_kwargs and not text_font_ok(s, **font_kwargs):
                continue
            return s
        return None

    def line_group(self, area, color_pred=None, width=None, arrow_count=0, min_count=1, horizontal=None, vertical=None, dashed=None):
        lines = self.line_candidates(area=area, color_pred=color_pred, width=width, horizontal=horizontal, vertical=vertical)
        if dashed is not None:
            if dashed:
                lines = [s for s in lines if is_dash_line(s) or min(shape_bounds(s)[2], shape_bounds(s)[3]) < 0.03]
            else:
                lines = [s for s in lines if not is_dash_line(s)]
        arrows = [s for s in lines if has_arrow_head(s)]
        return len(lines) >= min_count and len(arrows) >= arrow_count

    def bent_arrow_exists(self, area, color_pred=None, width=None, min_w=0.5, min_h=0.5):
        """检测单个折线/肘形箭头，而不是多段普通直线拼接。"""
        for s in self.line_candidates(area=area, color_pred=color_pred, width=width, arrow=True):
            l,t,w,h = shape_bounds(s)
            if w >= min_w and h >= min_h:
                return True
        return False

    def two_turn_path_in_region(self, region, dirs, color_pred=None, width=None, tol=0.4):
        """检测由三段直线首尾相连、拥有两个拐点的折线箭头。
        - region=(lx0,lx1,ty0,ty1)：整条折线各段都必须落在此范围（对应细则"距左/距上"）。
        - dirs=('right','down','right')：三段依次的方向（末段带箭头，即"最后指向"）。
        折线在办公软件中通常是多段独立直线段拼成（虚线更会拆成大量共线短段），
        先把共线相邻的短段合并成"整段"，再按端点相连 + 逐段方向 + 末段箭头判定。
        """
        def is_seg(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))

        lx0, lx1, ty0, ty1 = region
        raw = []
        for s in self.shapes:
            if not is_seg(s):
                continue
            if color_pred and not color_pred(get_line_color(s)):
                continue
            wpt = get_line_width_pt(s)
            if width and wpt is not None and not in_range(wpt, *width):
                continue
            l, t, w, h = shape_bounds(s)
            if not (lx0 - tol <= l and l + w <= lx1 + tol and ty0 - tol <= t and t + h <= ty1 + tol):
                continue
            raw.append(s)
        if not raw:
            return False

        # 把共线相邻的水平/竖直短段合并成整段：run=(orient,minx,miny,maxx,maxy,has_arrow)
        hsegs = [s for s in raw if is_horizontal(s) or (shape_bounds(s)[2] >= shape_bounds(s)[3])]
        vsegs = [s for s in raw if is_vertical(s) or (shape_bounds(s)[3] > shape_bounds(s)[2])]

        def merge(segments, orient):
            runs = []
            def bounds5(s):
                l, t, w, h = shape_bounds(s)
                return l, t, l + w, t + h, has_arrow_head(s)
            for s in segments:
                l, t, r, b, arr = bounds5(s)
                placed = False
                for run in runs:
                    if orient == 'H':
                        # 同一水平线(y接近) 且 x 方向相邻/重叠
                        if abs(((t + b) / 2) - ((run[2] + run[4]) / 2)) <= tol and \
                           l <= run[3] + tol and r >= run[1] - tol:
                            run[1] = min(run[1], l); run[3] = max(run[3], r)
                            run[2] = min(run[2], t); run[4] = max(run[4], b)
                            run[5] = run[5] or arr
                            placed = True; break
                    else:
                        if abs(((l + r) / 2) - ((run[1] + run[3]) / 2)) <= tol and \
                           t <= run[4] + tol and b >= run[2] - tol:
                            run[2] = min(run[2], t); run[4] = max(run[4], b)
                            run[1] = min(run[1], l); run[3] = max(run[3], r)
                            run[5] = run[5] or arr
                            placed = True; break
                if not placed:
                    runs.append([orient, l, t, r, b, arr])
            # 迭代合并直到稳定（相邻 run 可能因顺序未合并）
            changed = True
            while changed:
                changed = False
                for i in range(len(runs)):
                    for j in range(i + 1, len(runs)):
                        a, c = runs[i], runs[j]
                        if orient == 'H':
                            if abs(((a[2]+a[4])/2)-((c[2]+c[4])/2)) <= tol and \
                               a[1] <= c[3]+tol and a[3] >= c[1]-tol:
                                a[1]=min(a[1],c[1]);a[3]=max(a[3],c[3]);a[2]=min(a[2],c[2]);a[4]=max(a[4],c[4]);a[5]=a[5] or c[5]
                                runs.pop(j); changed=True; break
                        else:
                            if abs(((a[1]+a[3])/2)-((c[1]+c[3])/2)) <= tol and \
                               a[2] <= c[4]+tol and a[4] >= c[2]-tol:
                                a[2]=min(a[2],c[2]);a[4]=max(a[4],c[4]);a[1]=min(a[1],c[1]);a[3]=max(a[3],c[3]);a[5]=a[5] or c[5]
                                runs.pop(j); changed=True; break
                    if changed: break
            return runs

        runs = merge(hsegs, 'H') + merge(vsegs, 'V')
        if len(runs) < 3:
            return False

        def run_ends(run):
            orient, l, t, r, b, arr = run
            if orient == 'H':
                return (l, (t + b) / 2), (r, (t + b) / 2), arr
            return ((l + r) / 2, t), ((l + r) / 2, b), arr

        def near(p, q):
            return abs(p[0] - q[0]) <= tol and abs(p[1] - q[1]) <= tol

        def matches(want, p_from, p_to):
            if want in ('right', 'left'):
                if abs(p_from[1] - p_to[1]) > tol:
                    return False
                return (p_to[0] > p_from[0]) if want == 'right' else (p_to[0] < p_from[0])
            else:
                if abs(p_from[0] - p_to[0]) > tol:
                    return False
                return (p_to[1] > p_from[1]) if want == 'down' else (p_to[1] < p_from[1])

        # 枚举三段整段的有序组合，验证首尾相连 + 方向序列 + 末段带箭头
        for r1 in runs:
            a1, b1, _ = run_ends(r1)
            for start, end1 in ((a1, b1), (b1, a1)):
                if not matches(dirs[0], start, end1):
                    continue
                for r2 in runs:
                    if r2 is r1:
                        continue
                    a2, b2, _ = run_ends(r2)
                    for j2, e2 in ((a2, b2), (b2, a2)):
                        if not near(j2, end1) or not matches(dirs[1], j2, e2):
                            continue
                        for r3 in runs:
                            if r3 is r1 or r3 is r2:
                                continue
                            a3, b3, arr3 = run_ends(r3)
                            if not arr3:
                                continue
                            for j3, e3 in ((a3, b3), (b3, a3)):
                                if near(j3, e2) and matches(dirs[2], j3, e3):
                                    return True
        return False



    def _region_has_rounded_bent_arrow(self, region, color_pred=None, width=None,
                                       dashed_only=False, tol=0.4):
        """检测给定 region 内是否存在一个"单一形状"承载整条圆角折线箭头。
        "圆角"以办公软件实际渲染为准，命中以下之一：
          - `bentConnector*` / `curvedConnector*` 连接符（拐角默认圆角）；
          - `FREEFORM` 自由路径且含 `<a:arcTo>` 弧线段；
          - 折线形状且 `<a:ln>` 显式带 `<a:round/>` 连接样式。
        三段独立的 `prst="line"` 直线在同点相接（直角拐点）不算圆角。"""
        lx0, lx1, ty0, ty1 = region
        for s in self.shapes:
            if color_pred and not color_pred(get_line_color(s)):
                continue
            wpt = get_line_width_pt(s)
            if width and wpt is not None and not in_range(wpt, *width):
                continue
            l, t, w, h = shape_bounds(s)
            if not (lx0 - tol <= l and l + w <= lx1 + tol
                    and ty0 - tol <= t and t + h <= ty1 + tol):
                continue
            # 必须是"拥有拐角"的整段折线：跨度须同时具备水平和竖直方向
            if w < 0.4 or h < 0.4:
                continue
            if dashed_only and not is_dash_line(s):
                continue
            prst = ''
            try:
                g = s._element.spPr.find(qn('a:prstGeom'))
                if g is not None:
                    prst = g.get('prst', '') or ''
            except Exception:
                pass
            # ① bent/curved connector 连接符
            if prst.startswith('bentConnector') or prst.startswith('curvedConnector'):
                return True
            # ② freeform 且路径含 arcTo 弧段
            try:
                if s.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                    if s._element.spPr.find('.//' + qn('a:arcTo')) is not None:
                        return True
            except Exception:
                pass
            # ③ 折线形状且 line 显式 round join（排除单段 line 直线）
            try:
                ln = s._element.spPr.find(qn('a:ln'))
                if ln is not None and ln.find(qn('a:round')) is not None and prst != 'line':
                    return True
            except Exception:
                pass
        return False

    def two_turn_arrow_to_box(self, region, origin_box, target_box,
                              color_pred=None, width=None, vertical_dir='up', tol=0.4):
        """检测一条"先向右→再向上/向下→最后向右指向目标框"的两拐点折线箭头。
        - region=(lx0,lx1,ty0,ty1)：整条折线必须落在此范围内（对应细则"距左/距上"）。
        - origin_box：折线起点所依附的圆角矩形（文本所在框）。
        - target_box：箭头最终指向的圆角矩形。
        - vertical_dir：中间竖段方向，'up'（向上）或 'down'（向下）。
        折线在办公软件中可为多段相连的直线段拼成，本检测按端点相连+方向+箭头判定。
        """
        def is_seg(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))

        lx0, lx1, ty0, ty1 = region
        segs = []
        for s in self.shapes:
            if not is_seg(s):
                continue
            if color_pred and not color_pred(get_line_color(s)):
                continue
            wpt = get_line_width_pt(s)
            if width and wpt is not None and not in_range(wpt, *width):
                continue
            l, t, w, h = shape_bounds(s)
            if not (lx0 - tol <= l and l + w <= lx1 + tol and ty0 - tol <= t and t + h <= ty1 + tol):
                continue
            segs.append(s)
        if len(segs) < 3:
            return False

        hor = [s for s in segs if is_horizontal(s)]
        ver = [s for s in segs if is_vertical(s)]
        if not hor or not ver:
            return False

        def ends(s):
            l, t, w, h = shape_bounds(s)
            if is_horizontal(s):
                return (l, t + h / 2), (l + w, t + h / 2)  # 左端, 右端
            return (l + w / 2, t), (l + w / 2, t + h)      # 上端, 下端

        def near(p, q):
            return abs(p[0] - q[0]) <= tol and abs(p[1] - q[1]) <= tol

        obl, obt, obw, obh = shape_bounds(origin_box)
        tbl, tbt, tbw, tbh = shape_bounds(target_box)

        # 末段：水平、带箭头、右端指向目标框左缘
        for a in hor:
            if not has_arrow_head(a):
                continue
            a_left, a_right = ends(a)
            if abs(a_right[0] - tbl) > tol or not (tbt - tol <= a_right[1] <= tbt + tbh + tol):
                continue
            # 中间竖段：一端与末段左端相连
            for v in ver:
                v_top, v_bot = ends(v)
                if near(v_top, a_left):
                    join_a, other = v_top, v_bot     # a 在竖段顶端
                    a_at_top = True
                elif near(v_bot, a_left):
                    join_a, other = v_bot, v_top      # a 在竖段底端
                    a_at_top = False
                else:
                    continue
                # 方向校验：'up' 要求末段在竖段顶部，'down' 要求在底部
                if vertical_dir == 'up' and not a_at_top:
                    continue
                if vertical_dir == 'down' and a_at_top:
                    continue
                # 首段：水平，一端连竖段另一端，另一端贴近起点框右缘
                for h0 in hor:
                    if h0 is a:
                        continue
                    h_left, h_right = ends(h0)
                    if near(h_left, other) or near(h_right, other):
                        far = h_right if near(h_left, other) else h_left
                        if abs(far[0] - (obl + obw)) <= tol and (obt - tol <= far[1] <= obt + obh + tol):
                            return True
        return False


    def dashed_vertical_line_count(self, area, color_pred=None, width=None):
        """统计区域内真正的竖向虚线网格线。"""
        count = 0
        for s in self.line_candidates(area=area, color_pred=color_pred, width=width, vertical=True):
            if is_dash_line(s):
                count += 1
        return count

    def score(self):
        self.detect_key_objects()
        sh = self.shapes

        # 1. 第1页白色页面背景：覆盖整张幻灯片，颜色为白色或极浅灰白色，
        #    页面没有深色背景、水印。（逐点检查，且以办公软件真实生效的背景为准）
        page_area = self.slide_w * self.slide_h
        bg_kind, bg_rgb = self.background_kind_and_color()

        # 点1+点2：整页背景为白色或极浅灰白色纯色（'none' 时办公软件回落为白底）
        if bg_kind == 'none':
            bg_white_page = True
        elif bg_kind == 'solid':
            bg_white_page = color_is_white_or_grayish_white(bg_rgb)
        else:  # gradient / image / ref 等，非纯白整页背景
            bg_white_page = False

        # 点3：页面没有深色背景（背景本身不深，且没有铺满页面的深色形状充当背景）
        bg_self_dark = (bg_kind == 'solid' and bg_rgb is not None and not color_is_light(bg_rgb))
        dark_cover = [s for s in sh
                      if bbox_area(s) > page_area * 0.5
                      and not fill_is_none_or_transparent(s)
                      and not color_is_light(get_fill_color(s))]
        no_dark_bg = (not bg_self_dark) and (not dark_cover)

        # 点4：页面没有水印
        no_watermark = not self.has_watermark()

        self.add(3, '整页白色或极浅灰白色背景，覆盖整张幻灯片，无深色背景、无水印',
                 bg_white_page and no_dark_bg and no_watermark)

        # 2. 第1页左侧圆角矩形。逐点检查（位置/尺寸已在 detect_key_objects 的
        #    rrects 范围内约束：距左0.5–5.4cm、距上4.3–11.8cm、宽4.2–5.5cm、高5.5–7.2cm、
        #    且形状为圆角矩形）。此处补齐：填充浅青白色、边线青绿色单实线、线宽1–1.5磅。
        #    注：办公软件里常见两层圆角矩形叠合（底色壳 + 描边壳），detect_key_objects
        #    的 left_box 只取第一个，这里独立扫描全部位置尺寸匹配的圆角矩形候选，
        #    只要其中一个同时满足全部条件即视为达标。
        lb = self.key.get('left_box')
        left_ok = False
        lb_candidates = [s for s in self.shapes if looks_like_rect(s) and is_rounded_rect(s)
                         and in_range(emu2cm(s.left), 0.5, 5.4)
                         and in_range(emu2cm(s.top), 4.3, 11.8)
                         and in_range(abs(emu2cm(s.width)), 4.2, 5.5)
                         and in_range(abs(emu2cm(s.height)), 5.5, 7.2)]
        for cand in lb_candidates:
            lw = get_line_width_pt(cand)
            if (color_is_light_cyan(get_fill_color(cand))          # 填充为浅青白色
                    and color_is_cyan_green(get_line_color(cand))   # 边线为青绿色
                    and line_is_single_solid(cand)                  # 单实线
                    and lw is not None and in_range(lw, 1, 1.5)):   # 线宽1–1.5磅
                left_ok = True
                break
        self.add(3, '左侧圆角矩形：位置尺寸达标、浅青白填充、青绿色单实线边线、线宽1–1.5磅', left_ok)

        # 3. 第1页左侧"Multivariate Time Series"文本。逐点检查：
        #    ① 位于左侧圆角矩形的上半部分内；② 文本内容分两行排列；
        #    ③ 字体 Arial/Calibri；④ 字号10–18磅；⑤ 加粗；⑥ 颜色青绿色；⑦ 水平居中。
        mts_ok = False
        for s in self.find_text('Multivariate Time Series'):
            # ① 位于左侧圆角矩形上半部分内（文本中心落在 left_box 上半区）
            if lb is None:
                continue
            bl, bt, bw, bh = shape_bounds(lb)
            cx, cy = center(s)
            if not (bl <= cx <= bl + bw and bt <= cy <= bt + bh / 2):
                continue
            # ② 文本内容分两行排列
            if rendered_line_count(s) != 2:
                continue
            # ③④⑤⑥⑦ 字体/字号/加粗/颜色/水平居中（Arial 或 Calibri，10–18磅，加粗，青绿色，居中）
            if text_font_ok(s, fonts=('arial', 'calibri'), size=(10, 18), bold=True,
                            color_pred=color_is_cyan_green, center_align=True):
                mts_ok = True
                break
        self.add(3, '左侧圆角矩形上半部："Multivariate Time Series"两行、Arial/Calibri 10–18磅加粗青绿色居中', mts_ok)


        # 4. 第1页左侧折线图。逐点检查：
        #    ① 位于左侧圆角矩形的下半部分；② 含浅灰色竖向虚线网格；
        #    ③ 4条不同颜色的折线；④ 颜色包含蓝、橙、紫、绿；
        #    ⑤ 折线线宽0.5–2磅；⑥ 整体宽度2.7–4.2cm；⑦ 高度2.8–4.8cm。
        def _is_line_like(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))

        chart_ok = False
        if lb is not None:
            bl, bt, bw, bh = shape_bounds(lb)
            # ① 下半部分区域：left_box 下半区
            rx0, rx1 = bl, bl + bw
            ry0, ry1 = bt + bh / 2, bt + bh
            region = [s for s in self.shapes if _is_line_like(s)
                      and rx0 <= center(s)[0] <= rx1 and ry0 <= center(s)[1] <= ry1]

            # ② 浅灰色竖向虚线网格（青绿之外的浅灰竖线，且线型为虚线，构成网格需≥2条）
            grid = [s for s in region
                    if color_is_light_gray(get_line_color(s)) and is_vertical(s) and is_dash_line(s)]
            has_grid = len(grid) >= 2

            # ③④⑤ 4条折线：蓝/橙/紫/绿各至少1条，线宽0.5–2磅
            #     线宽未显式设置时（w=0/None）按 PowerPoint 默认线宽处理，
            #     避免把主题继承的线段误判为 0 磅
            def _colored(pred):
                return [s for s in region
                        if pred(get_line_color(s))
                        and line_width_in_range(s, 0.5, 2)]
            blue_ln = _colored(color_is_blue)
            orange_ln = _colored(color_is_orange)
            purple_ln = _colored(color_is_purple)
            green_ln = _colored(color_is_green)
            has_four_colors = all([blue_ln, orange_ln, purple_ln, green_ln])

            # ⑥⑦ 整体尺寸：折线图全部线元素的包围盒
            chart_shapes = grid + blue_ln + orange_ln + purple_ln + green_ln
            if chart_shapes:
                l = min(shape_bounds(s)[0] for s in chart_shapes)
                t = min(shape_bounds(s)[1] for s in chart_shapes)
                r = max(shape_bounds(s)[0] + shape_bounds(s)[2] for s in chart_shapes)
                b = max(shape_bounds(s)[1] + shape_bounds(s)[3] for s in chart_shapes)
                size_ok = in_range(r - l, 2.7, 4.2) and in_range(b - t, 2.8, 4.8)
            else:
                size_ok = False

            chart_ok = has_grid and has_four_colors and size_ok
        self.add(3, '左侧圆角矩形下半部折线图：浅灰竖向虚线网格＋蓝/橙/紫/绿4色折线(0.5–2磅)，整体2.7–4.2×2.8–4.8cm', chart_ok)

        # 5. 第1页左侧"Multivariate Time Series"文本所在圆角矩形，出现一条
        #    "先向右→再向上→最后指向右侧青色圆角虚线矩形"的折线箭头。逐点检查：
        #    ① 起点依附文本所在的 left_box；② 折线走向：右→上→右指向青色框；
        #    ③ 位于距左4.5–8.5cm、距上2.5–7.0cm范围内；④ 黑色；
        #    ⑤ 折线箭头（末端带箭头）；⑥ 拥有两个拐点；⑦ 线宽0.5–2磅。
        # 5. 第1页左侧"Multivariate Time Series"文本所在的圆角矩形，
        #    出现一条先向右→再向上→最后指向右侧青色圆角虚线矩形的折线箭头：
        #    位于距左4.5–8.5cm、距上2.5–7cm范围内，为黑色圆角折线箭头，拥有两个拐点；线宽0.5–2磅。
        #    逐点检查（细则每一个点都必须踩到，办公软件实际渲染为准；同一条箭头必须同时满足所有条件）：
        #    ① 起点=左侧"Multivariate Time Series"文本所在的圆角矩形（left_box）；
        #    ② 折线箭头（末端带箭头）；
        #    ③ 方向：先向右→再向上→最后向右；
        #    ④ 最后指向右侧"青色圆角虚线矩形"（seq_box）——该目标框须为圆角矩形、青绿色边线、虚线；
        #    ⑤ 距左4.5–8.5cm、距上2.5–7.0cm 范围内；
        #    ⑥ 黑色；
        #    ⑦ 圆角折线箭头（办公软件里的"圆角"呈现：bentConnector*/curvedConnector* 连接符、
        #       或 FREEFORM 路径含 a:arcTo 弧段、或线段显式 <a:round/> 连接样式）；
        #    ⑧ 两个拐点；⑨ 线宽0.5–2磅。
        cyan_arrow_ok = False
        origin_box = self.key.get('left_box')
        target_box = self.key.get('seq_box')
        if origin_box is not None and target_box is not None:
            # ①：origin_box 由 rrects 筛出，是圆角矩形（细则要求"文本所在的圆角矩形"，已满足）
            # ④：目标框须为圆角+青绿色边线+虚线
            target_qualifies = (is_rounded_rect(target_box)
                                and color_is_cyan_green(get_line_color(target_box))
                                and self.is_dashed_outline(target_box))
            if target_qualifies:
                obl, obt, obw, obh = shape_bounds(origin_box)
                tbl, tbt, tbw, tbh = shape_bounds(target_box)
                tol = 0.6
                lx0, lx1, ty0, ty1 = 4.5, 8.5, 2.5, 7.0
                for s in self.shapes:
                    # 线状形状（连接符/自由线/线形状）
                    if not (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                            or 'line' in shape_name_or_prst(s)
                            or 'connector' in shape_name_or_prst(s)):
                        continue
                    # ⑥ 黑色（含主题色 tx1 / 继承默认线条色的深色）
                    if not line_color_is_black_or_dark(s):
                        continue
                    # ⑨ 线宽0.5–2磅（含继承主题/默认线宽的形状）
                    if not line_width_in_range(s, 0.5, 2):
                        continue
                    # ⑤ 范围
                    l, t, w, h = shape_bounds(s)
                    if not (lx0 - 0.4 <= l and l + w <= lx1 + 0.4
                            and ty0 - 0.4 <= t and t + h <= ty1 + 0.4):
                        continue
                    # 取 prstGeom prst
                    prst = ''
                    try:
                        g = s._element.spPr.find(qn('a:prstGeom'))
                        if g is not None:
                            prst = g.get('prst', '') or ''
                    except Exception:
                        pass
                    # ⑦ "一体的折线箭头"：必须是单一形状本身承载整条两拐点折线，
                    #      即 bentConnector3/curvedConnector3 连接符、或 FREEFORM+arcTo 弧段、
                    #      或显式 <a:round/> 连接样式的折线形状。
                    #    普通 line/straightConnector1 只能画一条直线，天然不具备"两个拐点"，
                    #    也无法体现"先右→再上→最后右"的折线走向，即使多段拼接看起来像折线，
                    #    也不算"一体的折线箭头"——因此不再纳入判定，排除拼接线段冒充折线箭头。
                    is_rounded = False
                    if prst in ('bentConnector3', 'curvedConnector3'):
                        # 两个拐点：单一形状边界框需同时跨水平+竖直
                        if w >= 0.4 and h >= 0.4:
                            is_rounded = True
                    elif (prst.startswith('bentConnector')
                          or prst.startswith('curvedConnector')):
                        pass
                    else:
                        try:
                            if s.shape_type == MSO_SHAPE_TYPE.FREEFORM \
                                    and s._element.spPr.find('.//' + qn('a:arcTo')) is not None \
                                    and w >= 0.4 and h >= 0.4:
                                is_rounded = True
                        except Exception:
                            pass
                        if not is_rounded:
                            try:
                                ln = s._element.spPr.find(qn('a:ln'))
                                if ln is not None and ln.find(qn('a:round')) is not None \
                                        and w >= 0.4 and h >= 0.4:
                                    is_rounded = True
                            except Exception:
                                pass
                    if not is_rounded:
                        continue
                    # ② 折线箭头：末端带箭头
                    if not has_arrow_head(s):
                        continue
                    # ①③④ 圆角折线箭头几何走向：起点靠近 origin_box 右缘、终点靠近 target_box 左缘；
                    #      方向 right→up→right ⇒ 整体从 origin 中心 y 走到 target 中心 y（向上）；
                    #      箭头形状边界框：左缘≈origin 右缘、右缘≈target 左缘、
                    #      下缘覆盖 origin 中心 y、上缘覆盖 target 中心 y。
                    if not (abs(l - (obl + obw)) <= tol):
                        continue
                    if not (abs((l + w) - tbl) <= tol):
                        continue
                    origin_cy = obt + obh / 2
                    target_cy = tbt + tbh / 2
                    if not (target_cy < origin_cy):        # 目标框在起点框上方（"向上"）
                        continue
                    if not (t <= target_cy + tol and t + h >= origin_cy - tol):
                        continue
                    cyan_arrow_ok = True
                    break
        self.add(3, '左侧文本框到右侧青色圆角虚线矩形：黑色圆角折线箭头(先右→上→右、两拐点、0.5–2磅，距左4.5–8.5/距上2.5–7.0cm)', cyan_arrow_ok)

        # 6. 第1页左侧"Multivariate Time Series"文本所在的圆角矩形，
        #    出现一条先向右→再向下→最后向右指向紫色圆角虚线矩形的折线箭头：
        #    位于距左5.0–7.6cm、距上9.3–11.2cm范围内，为黑色圆角折线箭头，
        #    拥有两个拐点；线宽0.5–2磅。
        #    逐点检查（细则每一个点都必须踩到，办公软件实际渲染为准；同一条箭头必须同时满足所有条件）：
        #    ① 起点=左侧"Multivariate Time Series"文本所在的圆角矩形（left_box）；
        #    ② 折线箭头（末端带箭头）；
        #    ③ 方向：先向右→再向下→最后向右；
        #    ④ 最后指向右下"紫色圆角虚线矩形"（purple_box）——该目标框须为圆角矩形、紫色边线、虚线；
        #    ⑤ 距左5.0–7.6cm、距上9.3–11.2cm 范围内；
        #    ⑥ 黑色；
        #    ⑦ 圆角折线箭头（办公软件里的"圆角"呈现：bentConnector*/curvedConnector* 连接符、
        #       或 FREEFORM 路径含 a:arcTo 弧段、或线段显式 <a:round/> 连接样式）；
        #    ⑧ 两个拐点；⑨ 线宽0.5–2磅。
        purple_arrow_ok = False
        origin_box = self.key.get('left_box')
        purple_target = self.key.get('purple_box')
        if origin_box is not None and purple_target is not None:
            # ④ 目标框须为圆角+紫色边线+虚线
            target_qualifies = (is_rounded_rect(purple_target)
                                and color_is_purple(get_line_color(purple_target))
                                and self.is_dashed_outline(purple_target))
            if target_qualifies:
                obl, obt, obw, obh = shape_bounds(origin_box)
                tbl, tbt, tbw, tbh = shape_bounds(purple_target)
                tol = 0.6
                lx0, lx1, ty0, ty1 = 5.0, 7.6, 9.3, 11.2
                for s in self.shapes:
                    # 线状形状（连接符/自由线/线形状）
                    if not (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                            or 'line' in shape_name_or_prst(s)
                            or 'connector' in shape_name_or_prst(s)):
                        continue
                    # ⑥ 黑色（含主题色 tx1 / 继承默认线条色的深色）
                    if not line_color_is_black_or_dark(s):
                        continue
                    # ⑨ 线宽0.5–2磅（含继承主题/默认线宽的形状）
                    if not line_width_in_range(s, 0.5, 2):
                        continue
                    # ⑤ 范围
                    l, t, w, h = shape_bounds(s)
                    if not (lx0 - 0.4 <= l and l + w <= lx1 + 0.4
                            and ty0 - 0.4 <= t and t + h <= ty1 + 0.4):
                        continue
                    # 取 prstGeom prst
                    prst = ''
                    try:
                        g = s._element.spPr.find(qn('a:prstGeom'))
                        if g is not None:
                            prst = g.get('prst', '') or ''
                    except Exception:
                        pass
                    # ⑦⑧ "一体的折线箭头"：单一形状本身承载整条两拐点折线，
                    #     即 bentConnector3/curvedConnector3 连接符、或 FREEFORM+arcTo 弧段、
                    #     或显式 <a:round/> 连接样式的折线形状。
                    is_rounded = False
                    if prst in ('bentConnector3', 'curvedConnector3'):
                        # 两个拐点：单一形状边界框需同时跨水平+竖直
                        if w >= 0.4 and h >= 0.4:
                            is_rounded = True
                    elif (prst.startswith('bentConnector')
                          or prst.startswith('curvedConnector')):
                        pass
                    else:
                        try:
                            if s.shape_type == MSO_SHAPE_TYPE.FREEFORM \
                                    and s._element.spPr.find('.//' + qn('a:arcTo')) is not None \
                                    and w >= 0.4 and h >= 0.4:
                                is_rounded = True
                        except Exception:
                            pass
                        if not is_rounded:
                            try:
                                ln = s._element.spPr.find(qn('a:ln'))
                                if ln is not None and ln.find(qn('a:round')) is not None \
                                        and w >= 0.4 and h >= 0.4:
                                    is_rounded = True
                            except Exception:
                                pass
                    if not is_rounded:
                        continue
                    # ② 折线箭头：末端带箭头
                    if not has_arrow_head(s):
                        continue
                    # ①③④ 圆角折线箭头几何走向：起点靠近 origin_box 右缘、终点靠近 purple_box 左缘；
                    #      方向 right→down→right ⇒ 从 origin 中心 y 走到 target 中心 y（向下）；
                    #      箭头形状边界框：左缘≈origin 右缘、右缘≈target 左缘、
                    #      上缘覆盖 origin 中心 y、下缘覆盖 target 中心 y。
                    if not (abs(l - (obl + obw)) <= tol):
                        continue
                    if not (abs((l + w) - tbl) <= tol):
                        continue
                    origin_cy = obt + obh / 2
                    target_cy = tbt + tbh / 2
                    if not (target_cy > origin_cy):        # 目标框在起点框下方（"向下"）
                        continue
                    if not (t <= origin_cy + tol and t + h >= target_cy - tol):
                        continue
                    purple_arrow_ok = True
                    break
        self.add(3, '左侧文本框到右下紫色圆角虚线矩形：黑色圆角折线箭头(先右→下→右、两拐点、0.5–2磅，距左5.0–7.6/距上9.3–11.2cm)', purple_arrow_ok)

        # 7. 第1页左上方矩形形状。逐点检查：① 位于距左7.8–14.4cm、距上1.2–7.0cm
        #    （已由 detect_key_objects 的 rrects 约束）；② 圆角矩形；③ 青绿色虚线边框；
        #    ④ 宽5.5–6.7cm、高3.6–5.6cm（已由 rrects 约束）；⑤ 线宽0.5–2磅；⑥ 内部无填充。
        seq = self.key.get('seq_box')
        # 7. 第1页左上方出现矩形形状：位于距左7.8–14.4cm、距上1.2–7.0cm范围内，
        #    形状为青绿色虚线圆角矩形，宽5.5–6.7cm，高3.6–5.6cm，线宽0.5–2磅；内部无填充。
        #    逐点检查（细则每一个点都必须踩到，办公软件实际渲染为准）：
        #    ① 距左7.8–14.4cm；② 距上1.2–7.0cm；③ 圆角矩形；④ 青绿色边线；⑤ 虚线；
        #    ⑥ 宽5.5–6.7cm；⑦ 高3.6–5.6cm；⑧ 线宽0.5–2磅；⑨ 内部无填充。
        seq = None
        for s in self.shapes:
            if not looks_like_rect(s):                                # ③ 前提：矩形
                continue
            if not is_rounded_rect(s):                                # ③ 圆角矩形
                continue
            if not box_in(s, left=(7.8, 14.4), top=(1.2, 7.0),         # ①② 位置
                          width=(5.5, 6.7), height=(3.6, 5.6)):        # ⑥⑦ 尺寸
                continue
            if not color_is_cyan_green(get_line_color(s)):            # ④ 青绿色边线
                continue
            if not self.is_dashed_outline(s):                         # ⑤ 虚线
                continue
            lw = get_line_width_pt(s)
            if not (lw is not None and in_range(lw, 0.5, 2)):         # ⑧ 线宽0.5–2磅
                continue
            if not fill_is_visually_empty(s):                         # ⑨ 内部无填充
                continue
            seq = s
            break
        # 更新 key，便于后续依赖此框的检查（如 Sequence Encoder 文本、GRU 矩形等）
        self.key['seq_box'] = seq
        seq_ok = seq is not None
        self.add(5, '左上青绿色虚线圆角矩形：距左7.8–14.4/距上1.2–7.0cm、宽5.5–6.7/高3.6–5.6cm、线宽0.5–2磅、内部无填充', seq_ok)

        # 8. 第1页上方"Sequence Encoder"文本。逐点检查：
        #    ① 位于青色虚线圆角矩形上方；② 文本内容仅排列为1行；
        #    ③ 字体 Arial 或 Calibri；④ 字号12–18磅；⑤ 加粗；⑥ 颜色为青绿色。
        seq_text_ok = False
        if seq is not None:
            bl, bt, bw, bh = shape_bounds(seq)
            for s in self.find_text('Sequence Encoder'):
                # ② 仅1行
                if rendered_line_count(s) != 1:
                    continue
                # ① 位于青色框上方：文本中心在框顶边以上，且横向与框重叠（属于该框）
                cx, cy = center(s)
                tl, _tt, tw, _th = shape_bounds(s)
                horiz_overlap = min(tl + tw, bl + bw) - max(tl, bl) > 0
                if not (cy <= bt + 0.1 and horiz_overlap):
                    continue
                # ③④⑤⑥ 字体 Arial/Calibri、字号12–18磅、加粗、青绿色
                if text_font_ok(s, fonts=('arial', 'calibri'), size=(12, 18), bold=True,
                                color_pred=color_is_cyan_green):
                    seq_text_ok = True
                    break
        self.add(3, '青色框上方"Sequence Encoder"：仅1行、Arial/Calibri 12–18磅加粗青绿色', seq_text_ok)

        # 9. 第1页"Sequence Encoder"下方两个圆角矩形。逐点检查：
        #    ① 位于青色虚线圆角矩形内；② 浅青色填充；③ 实线（非虚线）；④ 圆角矩形；
        #    ⑤ 边线青绿色单实线；⑥ 线宽0.5–2.0磅；⑦ 两个矩形内部均含"Gated Recurrent Unit"文本；
        #    ⑧ 文本分两行排列；⑨ 字号9–18磅；⑩ 加粗；⑪ 颜色黑色。
        good_gru = []
        if seq is not None:
            bl, bt, bw, bh = shape_bounds(seq)
            for txt in self.find_text('Gated Recurrent Unit'):
                # ⑧⑨⑩⑪ 文本：两行、字号9–18磅、加粗、黑色
                if rendered_line_count(txt) != 2:
                    continue
                if not text_font_ok(txt, size=(9, 18), bold=True, color_pred=color_is_black_or_dark):
                    continue
                # ②③④⑤⑥ 承载文本的圆角矩形：浅青填充、实线、青绿色单实线、线宽0.5–2磅
                def _rect_pred(r):
                    lw = get_line_width_pt(r)
                    return (color_is_light_cyan(get_fill_color(r))       # 浅青色填充
                            and not self.is_dashed_outline(r)             # 实线（非虚线）
                            and color_is_cyan_green(get_line_color(r))    # 青绿色边线
                            and line_is_single_solid(r)                   # 单实线
                            and lw is not None and in_range(lw, 0.5, 2))  # 线宽0.5–2.0磅
                rect = self.rect_covering_text(txt, _rect_pred)           # ④ 圆角矩形（函数只取圆角矩形）
                if rect is None:
                    continue
                # ① 位于青色虚线圆角矩形内
                cx, cy = center(rect)
                if bl <= cx <= bl + bw and bt <= cy <= bt + bh and rect not in good_gru:
                    good_gru.append(rect)
        gru_ok = len(good_gru) >= 2
        self.add(5, '青色框内两个浅青实线圆角矩形(青绿单实线0.5–2磅)，各含两行9–18磅加粗黑色"Gated Recurrent Unit"', gru_ok)

        # 10. 第1页"Gated Recurrent Unit"文本右侧折线箭头。逐点检查（细则每一个点都必须踩到，
        #    办公软件实际渲染为准；同一条箭头必须同时满足所有条件）：
        #    ① 位于距左12.9–15.5cm、距上2.8–5.5cm范围内；
        #    ② 黑色；
        #    ③ 线宽0.5–2磅；
        #    ④ 折线箭头（末端带箭头）；
        #    ⑤ 圆角折线箭头（办公软件里"圆角"呈现：bentConnector3/curvedConnector3 连接符、
        #       或 FREEFORM 路径含 a:arcTo 弧段、或折线形状显式带 <a:round/> 连接样式）；
        #    ⑥ 两个拐点（bentConnector3/curvedConnector3 天然=2 拐点；三段独立 prst="line"
        #       直角相接不算圆角折线，此处不予认可）；
        #    ⑦ 方向：先向右→再向下→最后指向右侧（起点靠区域左上、终点在右侧且位于起点下方）。
        gru_arrow_ok = False
        lx0, lx1, ty0, ty1 = 12.9, 15.5, 2.8, 5.5
        tol = 0.4
        for s in self.shapes:
            if not (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s)
                    or 'connector' in shape_name_or_prst(s)):
                continue
            # ② 黑色（含主题色 tx1 / 继承默认线条色的深色）
            if not line_color_is_black_or_dark(s):
                continue
            # ③ 线宽0.5–2磅（含继承主题/默认线宽）
            if not line_width_in_range(s, 0.5, 2):
                continue
            # ① 位置
            l, t, w, h = shape_bounds(s)
            if not (lx0 - tol <= l and l + w <= lx1 + tol
                    and ty0 - tol <= t and t + h <= ty1 + tol):
                continue
            # ⑤ "一体的折线箭头"：必须是单一形状本身承载整条两拐点折线，
            #      即 bentConnector3/curvedConnector3 连接符、或 FREEFORM+arcTo 弧段、
            #      或显式 <a:round/> 连接样式的折线形状。
            #    普通 line/straightConnector1 只能画一条直线，天然不具备"两个拐点"，
            #    也无法体现"先右→再下→最后右"的折线走向，多段拼接冒充折线箭头不予认可。
            prst = ''
            try:
                g = s._element.spPr.find(qn('a:prstGeom'))
                if g is not None:
                    prst = g.get('prst', '') or ''
            except Exception:
                pass
            is_rounded = False
            if prst in ('bentConnector3', 'curvedConnector3'):
                if w >= 0.4 and h >= 0.4:
                    is_rounded = True
            elif (prst.startswith('bentConnector')
                  or prst.startswith('curvedConnector')):
                # bentConnector2=1拐点、bentConnector4=3拐点 —— 不满足"两拐点"
                pass
            else:
                try:
                    if s.shape_type == MSO_SHAPE_TYPE.FREEFORM \
                            and s._element.spPr.find('.//' + qn('a:arcTo')) is not None \
                            and w >= 0.4 and h >= 0.4:
                        is_rounded = True
                except Exception:
                    pass
                if not is_rounded:
                    try:
                        ln = s._element.spPr.find(qn('a:ln'))
                        if ln is not None and ln.find(qn('a:round')) is not None \
                                and w >= 0.4 and h >= 0.4:
                            is_rounded = True
                    except Exception:
                        pass
            if not is_rounded:
                continue
            # ④ 折线箭头：末端带箭头
            if not has_arrow_head(s):
                continue
            xfrm = s._element.find('.//' + qn('a:xfrm'))
            flip_h = xfrm is not None and xfrm.get('flipH') == '1'
            flip_v = xfrm is not None and xfrm.get('flipV') == '1'
            # ⑦ 方向"右→下→右"：末端在右下、起点在左上；
            #    箭头形状为单一折线，其 bbox 的右下角=末端(右下)，左上角=起点(左上)。
            #    右→下→右 意味着终点 y 大于起点 y，且终点 x 大于起点 x（整体向右下延伸）。
            #    对 bent/curved connector，flip 属性会改变实际走向，这里以 bbox+flip 综合判定：
            #    - 未 flipH 且未 flipV：起点在左上，末端在右下 → 满足"右→下→右"
            if flip_h or flip_v:
                # 走向被翻转，不再是"右→下→右"
                continue
            gru_arrow_ok = True
            break
        self.add(3, 'GRU右侧黑色圆角折线箭头：先右→下→右、两拐点、线宽0.5–2磅(距左12.9–15.5/距上2.8–5.5cm)', gru_arrow_ok)

        # 11. 第1页左侧紫色矩形形状。逐点检查：① 位于距左7.2–13.8cm、距上7.5–12.8cm
        #    （已由 detect_key_objects 的 rrects 约束）；② 圆角矩形；③ 虚线；④ 紫色边线；
        #    ⑤ 宽5–6cm、高4.7–5.6cm（已由 rrects 约束）；⑥ 线宽0.5–2磅；
        #    ⑦ 填充为无填充或淡紫色（以办公软件实际渲染为准：近白视为无填充）。
        pur = self.key.get('purple_box')
        pur_ok = False
        if pur is not None:
            lw = get_line_width_pt(pur)
            pur_ok = (is_rounded_rect(pur)                              # 圆角矩形
                      and self.is_dashed_outline(pur)                   # 虚线
                      and color_is_purple(get_line_color(pur))          # 紫色边线
                      and lw is not None and in_range(lw, 0.5, 2)       # 线宽0.5–2磅
                      and (fill_is_visually_empty(pur)                  # 无填充（含近白视觉无填充）
                           or color_is_light_purple(get_fill_color(pur))))  # 或淡紫色
        self.add(5, '左侧紫色虚线圆角矩形：位置尺寸达标、线宽0.5–2磅、无填充或淡紫色填充', pur_ok)

        # 12. 第1页左侧紫色虚线矩形内部"2. Pattern Encoder"文本。逐点检查：
        #    ① 位于紫色矩形内部；② 位于矩形顶部；③ 文本仅1行；
        #    ④ 字体 Arial 或 Calibri；⑤ 字号12–18磅；⑥ 加粗；⑦ 字体颜色为紫色。
        pat_ok = False
        if pur is not None:
            pl, pt_, pw, ph = shape_bounds(pur)
            for s in self.find_text('2. Pattern Encoder'):
                # ③ 仅1行
                if rendered_line_count(s) != 1:
                    continue
                sl, st, sw, shh = shape_bounds(s)
                cx, cy = center(s)
                # ① 位于紫色矩形内部（文本中心落在框内）
                inside = pl <= cx <= pl + pw and pt_ <= cy <= pt_ + ph
                # ② 位于矩形顶部（文本中心在矩形上半部）
                at_top = cy <= pt_ + ph / 2
                if not (inside and at_top):
                    continue
                # ④⑤⑥⑦ 字体 Arial/Calibri、字号12–18磅、加粗、紫色
                if text_font_ok(s, fonts=('arial', 'calibri'), size=(12, 18), bold=True,
                                color_pred=color_is_purple):
                    pat_ok = True
                    break
        self.add(3, '紫色框内顶部"2. Pattern Encoder"：仅1行、Arial/Calibri 12–18磅加粗紫色', pat_ok)

        # 13. 第1页"Gated Recurrent Unit"下方矩形组合图块（方格矩阵）。逐点检查：
        #    ① 位于紫色虚线矩形中部，距左7.8–13.8cm、距上8.7–12.5cm；② 含紫色、浅紫色、白色小方格；
        #    ③ 矩阵5列×3行；④ 边线为紫色；⑤ 线宽0.5–1磅。
        cells = []
        for s in sh:
            if not looks_like_rect(s):
                continue
            if not box_in(s, left=(7.8, 13.8), top=(8.7, 12.5), width=(0.15, 1.5), height=(0.15, 1.5)):
                continue
            lw = get_line_width_pt(s)
            # ④ 紫色边线；⑤ 线宽0.5–1磅
            if not (color_is_purple(get_line_color(s)) and lw is not None and in_range(lw, 0.5, 1)):
                continue
            cells.append(s)

        # ③ 5列×3行：按左/上坐标聚类出的列数、行数
        def _cluster(vals, tol=0.25):
            groups = []
            for v in sorted(vals):
                if not groups or v - groups[-1][-1] > tol:
                    groups.append([v])
                else:
                    groups[-1].append(v)
            return len(groups)
        cols = _cluster([shape_bounds(s)[0] for s in cells])
        rows = _cluster([shape_bounds(s)[1] for s in cells])
        grid_ok = (cols == 5 and rows == 3 and len(cells) >= 15)

        # ② 三种填充：紫色、浅紫色、白色均出现
        has_purple = any(color_is_purple(get_fill_color(s)) for s in cells)
        has_light_purple = any(color_is_light_purple(get_fill_color(s)) for s in cells)
        has_white = any(color_is_white_or_near(get_fill_color(s)) for s in cells)
        colors_ok = has_purple and has_light_purple and has_white

        self.add(3, '紫色框中部矩形组合图块：5列×3行方格、紫色细线(0.5–1磅)、含紫/浅紫/白三色', grid_ok and colors_ok)

        # 14. 第1页"Temporal Embedding"文本。逐点检查：
        #    ① 位于紫色方格矩阵下方；② 文本分两行排列；③ 字体 Arial 或 Calibri；
        #    ④ 字号9–18磅；⑤ 颜色黑色；⑥ 加粗；⑦ 水平居中。
        # 矩阵底边（cells 来自第13项，取其最大下边界作为"下方"的参照）
        matrix_bottom = max((shape_bounds(s)[1] + shape_bounds(s)[3] for s in cells), default=None)
        temp_ok = False
        for s in self.find_text('Temporal Embedding'):
            # ② 两行
            if rendered_line_count(s) != 2:
                continue
            # ① 位于矩阵下方：文本中心在矩阵底边之下
            _cx, cy = center(s)
            if matrix_bottom is not None and cy <= matrix_bottom:
                continue
            # ③④⑤⑥⑦ Arial/Calibri、9–18磅、黑色、加粗、水平居中
            if text_font_ok(s, fonts=('arial', 'calibri'), size=(9, 18), bold=True,
                            color_pred=color_is_black_or_dark, center_align=True):
                temp_ok = True
                break
        self.add(3, '矩阵下方"Temporal Embedding"：两行、Arial/Calibri 9–18磅加粗黑色居中', temp_ok)

        # 15. 第1页紫色虚线矩形右侧折线箭头。逐点检查（细则每一个点都必须踩到，
        #    办公软件实际渲染为准；同一条箭头必须同时满足所有条件）：
        #    ① 位于距左12.7–15.5cm、距上6.8–10.3cm范围内；
        #    ② 紫色；
        #    ③ 线宽0.5–2磅；
        #    ④ 折线箭头（末端带箭头）；
        #    ⑤ 虚线（形状自身 <a:prstDash> 为真正的虚线；不接受"透明实线 + 拼装短线段"）；
        #    ⑥ 圆角箭头（办公软件里"圆角"呈现：bentConnector3/curvedConnector3 连接符、
        #       或 FREEFORM 路径含 a:arcTo 弧段、或折线形状显式带 <a:round/> 连接样式）；
        #    ⑦ 两个拐点（bentConnector3/curvedConnector3 天然=2 拐点；bentConnector2/4 排除）；
        #    ⑧ 方向：先向右→再向上→最后指向右侧（起点靠区域左下、终点在右上）。
        pur_arrow_ok = False
        lx0, lx1, ty0, ty1 = 12.7, 15.5, 6.8, 10.3
        tol = 0.4
        for s in self.shapes:
            if not (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s)
                    or 'connector' in shape_name_or_prst(s)):
                continue
            # ② 紫色
            if not color_is_purple(get_line_color(s)):
                continue
            # ③ 线宽0.5–2磅（含继承主题/默认线宽）
            if not line_width_in_range(s, 0.5, 2):
                continue
            # ① 位置
            l, t, w, h = shape_bounds(s)
            if not (lx0 - tol <= l and l + w <= lx1 + tol
                    and ty0 - tol <= t and t + h <= ty1 + tol):
                continue
            # ⑤ 虚线：形状自身 prstDash 为真正的虚线
            if not is_dash_line(s):
                continue
            # ⑥ "一体的折线箭头"：必须是单一形状本身承载整条两拐点折线，
            #      即 bentConnector3/curvedConnector3 连接符、或 FREEFORM+arcTo 弧段、
            #      或显式 <a:round/> 连接样式的折线形状。
            #    普通 line/straightConnector1 只能画一条直线，天然不具备"两个拐点"，
            #    也无法体现"先右→再上→最后右"的折线走向，多段拼接冒充折线箭头不予认可。
            prst = ''
            try:
                g = s._element.spPr.find(qn('a:prstGeom'))
                if g is not None:
                    prst = g.get('prst', '') or ''
            except Exception:
                pass
            is_rounded = False
            if prst in ('bentConnector3', 'curvedConnector3'):
                if w >= 0.4 and h >= 0.4:
                    is_rounded = True
            elif (prst.startswith('bentConnector')
                  or prst.startswith('curvedConnector')):
                pass
            else:
                try:
                    if s.shape_type == MSO_SHAPE_TYPE.FREEFORM \
                            and s._element.spPr.find('.//' + qn('a:arcTo')) is not None \
                            and w >= 0.4 and h >= 0.4:
                        is_rounded = True
                except Exception:
                    pass
                if not is_rounded:
                    try:
                        ln = s._element.spPr.find(qn('a:ln'))
                        if ln is not None and ln.find(qn('a:round')) is not None \
                                and w >= 0.4 and h >= 0.4:
                            is_rounded = True
                    except Exception:
                        pass
            if not is_rounded:
                continue
            # ④ 折线箭头：末端带箭头
            if not has_arrow_head(s):
                continue
            # ⑧ 方向"右→上→右"：末端在右上、起点在左下；
            #    对 bent/curved connector，flip 属性会改变实际走向，这里以 flip 综合判定。
            #    未 flipH 未 flipV：起点在 bbox 左下、末端在右上 → 满足"右→上→右"。
            xfrm = s._element.find('.//' + qn('a:xfrm'))
            flip_h = xfrm is not None and xfrm.get('flipH') == '1'
            flip_v = xfrm is not None and xfrm.get('flipV') == '1'
            # "右→上→右"要求整体向右+向上：意味着相对未 flip 的 bentConnector3
            # (起点左上→终点右下) 需要 flipV=1；flipH 不能变（否则末段变向左）。
            if flip_h:
                continue
            if not flip_v:
                continue
            pur_arrow_ok = True
            break
        self.add(3, '紫色框右侧紫色虚线圆角折线箭头：先右→上→右、两拐点、线宽0.5–2磅(距左12.7–15.5/距上6.8–10.3cm)', pur_arrow_ok)

        # 16. 第1页橙色矩形形状。逐点检查：① 位于距左15.1–29.7cm、距上1.2–9.3cm
        #    （已由 detect_key_objects 的 rrects 约束）；② 圆角矩形；③ 虚线；④ 橙色边线；
        #    ⑤ 宽13–15cm、高6–8cm（已由 rrects 约束）；⑥ 线宽0.5–2磅；
        #    ⑦ 内部无填充或填充为浅橙色（以办公软件实际渲染为准：近白视为无填充）。
        org = self.key.get('orange_box')
        org_ok = False
        if org is not None:
            lw = get_line_width_pt(org)
            org_ok = (is_rounded_rect(org)                             # 圆角矩形
                      and self.is_dashed_outline(org)                  # 虚线
                      and color_is_orange(get_line_color(org))         # 橙色边线
                      and lw is not None and in_range(lw, 0.5, 2)      # 线宽0.5–2磅
                      and (fill_is_visually_empty(org)                 # 无填充（含近白视觉无填充）
                           or color_is_light_orange(get_fill_color(org))))  # 或浅橙色
        self.add(5, '橙色虚线圆角矩形：位置尺寸达标、线宽0.5–2磅、无填充或浅橙色填充', org_ok)

        # 17. 第1页"3. Fusion & Representation Learning"文本。逐点检查：
        #    ① 位于橙色虚线矩形顶部上方；② 文本仅1行；③ 字体 Arial 或 Calibri；
        #    ④ 字号12–18磅；⑤ 加粗；⑥ 颜色橙色；⑦ 水平居中。
        fusion_ok = False
        if org is not None:
            ol, ot, ow, oh = shape_bounds(org)
            for s in self.find_text('3. Fusion & Representation Learning'):
                # ② 仅1行
                if rendered_line_count(s) != 1:
                    continue
                sl, st, sw, sh2 = shape_bounds(s)
                cx, cy = center(s)
                # ① 位于橙色框顶部上方：文本中心在框顶边以上，且横向与框重叠
                horiz_overlap = min(sl + sw, ol + ow) - max(sl, ol) > 0
                if not (cy <= ot + 0.1 and horiz_overlap):
                    continue
                # ③④⑤⑥⑦ Arial/Calibri、12–18磅、加粗、橙色、水平居中
                if text_font_ok(s, fonts=('arial', 'calibri'), size=(12, 18), bold=True,
                                color_pred=color_is_orange, center_align=True):
                    fusion_ok = True
                    break
        self.add(3, '橙色框顶部上方"3. Fusion & Representation Learning"：仅1行、Arial/Calibri 12–18磅加粗橙色居中', fusion_ok)

        # 18. 第1页橙色虚线矩形中4个大小一致的矩形形状。逐点检查：
        #    ① 橙框内有4个矩形；② 大小一致；③ 位于同一水平线上；④ 间距相等；
        #    ⑤ 圆角矩形；⑥ 浅橙色填充；⑦ 边线单实线；⑧ 线宽0.5–2磅；
        #    ⑨ 4个矩形从左至右顶部依次出现文本 Concatenate/Self-Attention/Feed-Forward/Layer Normalization。
        module_names = ['Concatenate', 'Self-Attention', 'Feed-Forward', 'Layer Normalization']
        modules_ok = False
        modules = []
        if org is not None:
            ol, ot, ow, oh = shape_bounds(org)
            cand = []
            for s in sh:
                if s is org or not (looks_like_rect(s) and is_rounded_rect(s)):  # ⑤ 圆角矩形
                    continue
                cx, cy = center(s)
                if not (ol <= cx <= ol + ow and ot <= cy <= ot + oh):           # ① 橙框内
                    continue
                lw = get_line_width_pt(s)
                if not (color_is_orangey_light_fill(get_fill_color(s))          # ⑥ 浅橙色填充
                        and line_is_single_solid(s)                             # ⑦ 单实线
                        and color_is_orange(get_line_color(s))                  # 边线橙色
                        and lw is not None and in_range(lw, 0.5, 2)):           # ⑧ 线宽0.5–2磅
                    continue
                cand.append(s)
            modules = sorted(cand, key=lambda s: center(s)[0])
            if len(modules) >= 4:
                modules = modules[:4]
                ws = [shape_bounds(s)[2] for s in modules]
                hs = [shape_bounds(s)[3] for s in modules]
                cys = [center(s)[1] for s in modules]
                cxs = [center(s)[0] for s in modules]
                # ② 大小一致（宽高相差在 0.3cm 内）
                size_same = (max(ws) - min(ws) <= 0.3) and (max(hs) - min(hs) <= 0.3)
                # ③ 同一水平线（中心 y 相差 <0.5cm）
                same_line = max(cys) - min(cys) < 0.5
                # ④ 间距相等（相邻中心间距相差在 0.3cm 内）
                gaps = [cxs[i + 1] - cxs[i] for i in range(3)]
                equal_gap = (max(gaps) - min(gaps)) <= 0.3
                # ⑨ 从左至右每个矩形顶部依次出现对应文本
                order_text_ok = True
                for rect, name in zip(modules, module_names):
                    rl, rt, rw, rh = shape_bounds(rect)
                    hit = False
                    for s in self.find_text(name):
                        tcx, tcy = center(s)
                        if rl <= tcx <= rl + rw and rt - 0.3 <= tcy <= rt + rh / 2:
                            hit = True
                            break
                    if not hit:
                        order_text_ok = False
                        break
                modules_ok = size_same and same_line and equal_gap and order_text_ok
        self.add(5, '橙框内4个大小一致、等间距、同水平线的浅橙圆角矩形(单实线0.5–2磅)，顶部依次为4个模块标题', modules_ok)

        # 19. 4个模块文本的字体/字号/颜色/行数。逐点检查：
        #    ① 4个文本字体为 Arial 或 Calibri；② 字号10–15磅；③ 颜色黑色；
        #    ④ Concatenate 仅1行；⑤ 其余3个（Self-Attention/Feed-Forward/
        #       Layer Normalization）各分两行（要求文本框可见文本去掉空白后精确等于给定名字，
        #       不接受额外内容）。
        def _text_equals(shape, name):
            raw = ''.join(shape_text(shape).split())
            target = ''.join(name.split())
            return raw == target
        text_ok = True
        for name in module_names:
            matched = [s for s in self.shapes if _text_equals(s, name)]
            if not matched:
                text_ok = False
                break
            s = matched[0]
            # ①②③ Arial/Calibri、10–15磅、黑色
            if not text_font_ok(s, fonts=('arial', 'calibri'), size=(10, 15),
                                color_pred=color_is_black_or_dark):
                text_ok = False
                break
            # ④⑤ 行数：Concatenate=1，其余=2
            want_lines = 1 if name == 'Concatenate' else 2
            if rendered_line_count(s) != want_lines:
                text_ok = False
                break
        self.add(3, '4个模块文本：Arial/Calibri 10–15磅黑色，Concatenate 1行、其余各名字2行', text_ok)

        # 20. 第1页"Concatenate"文本下方2列矩形组合阵列。逐点检查：
        #    ① 位于橙色矩形内部下方（Concatenate 文本下方）；② 8个橙色小方格矩形；
        #    ③ 分两列排布；④ 两列不重合且留有间距；⑤ 间距之间出现4条直线相连，线宽0.5–1磅；
        #    ⑥ 阵列整体宽1.5–2.5cm、高1.6–3.3cm。
        def _is_line_like(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))

        concat_ok = False
        cc = self.find_text('Concatenate')
        if org is not None and cc:
            ol, ot, ow, oh = shape_bounds(org)
            cl, ct, cw, ch = shape_bounds(cc[0])
            cc_cx = cl + cw / 2
            below_top = ct + ch  # Concatenate 文本底边，阵列在其下方
            # ② 橙色小方格：橙色（填充或边线）、位于橙框内、在文本下方、且横向靠近 Concatenate 列
            cells = [s for s in sh
                     if looks_like_rect(s) and s is not org
                     and (color_is_orange(get_line_color(s)) or color_is_orange(get_fill_color(s))
                          or color_is_light_orange(get_fill_color(s)))
                     and ol <= center(s)[0] <= ol + ow and ot <= center(s)[1] <= ot + oh
                     and center(s)[1] >= below_top
                     and abs(center(s)[0] - cc_cx) <= ow / 3
                     and shape_bounds(s)[2] < 1.3 and shape_bounds(s)[3] < 1.3]
            if len(cells) == 8:
                # ③ 两列：按中心 x 聚成的列数
                def _cluster(vals, tol=0.25):
                    groups = []
                    for v in sorted(vals):
                        if not groups or v - groups[-1][-1] > tol:
                            groups.append([v])
                        else:
                            groups[-1].append(v)
                    return groups
                col_groups = _cluster([center(s)[0] for s in cells])
                two_cols = len(col_groups) == 2
                gap_ok = False
                conn_ok = False
                if two_cols:
                    left_col = [s for s in cells if center(s)[0] in col_groups[0] or round(center(s)[0], 4) <= col_groups[0][-1] + 0.01]
                    # ④ 两列不重合且留有间距：左列右缘 < 右列左缘
                    left_right = max(shape_bounds(s)[0] + shape_bounds(s)[2]
                                     for s in cells if center(s)[0] <= col_groups[0][-1] + 0.01)
                    right_left = min(shape_bounds(s)[0]
                                     for s in cells if center(s)[0] >= col_groups[1][0] - 0.01)
                    gap_ok = right_left > left_right + 0.02
                    # ⑤ 间距之间4条直线相连、线宽0.5–1磅
                    conn = [s for s in sh if _is_line_like(s)
                            and get_line_width_pt(s) is not None and in_range(get_line_width_pt(s), 0.5, 1)
                            and left_right - 0.1 <= center(s)[0] <= right_left + 0.1]
                    conn_ok = len(conn) >= 4
                # ⑥ 阵列整体尺寸
                l = min(shape_bounds(s)[0] for s in cells)
                r = max(shape_bounds(s)[0] + shape_bounds(s)[2] for s in cells)
                t = min(shape_bounds(s)[1] for s in cells)
                b = max(shape_bounds(s)[1] + shape_bounds(s)[3] for s in cells)
                size_ok = in_range(r - l, 1.5, 2.5) and in_range(b - t, 1.6, 3.3)
                concat_ok = two_cols and gap_ok and conn_ok and size_ok
        self.add(3, 'Concatenate下方阵列：橙框内8个橙色小方格分两列有间距、间距间4条连线(0.5–1磅)、整体1.5–2.5×1.6–3.3cm', concat_ok)

        # 21. 第1页"Self-Attention"文本下方组合形状。逐点检查：
        #    ① 位于 Self-Attention 文本下方；② 橙色节点圆点；③ 竖向线条；④ 底部小矩形；
        #    ⑤ 三条虚线连接线；⑥ 三条虚线箭头；⑦ 整体宽1.8–2.4cm、高1.7–2.8cm；
        #    ⑧ 线条为橙色；⑨ 线宽0.5–1.25磅。
        def _is_line_like2(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))
        def _is_dot(s):
            n = shape_name_or_prst(s)
            return 'ellipse' in n or 'oval' in n or 'circle' in n

        att_ok = False
        sa = self.find_text('Self-Attention')
        if org is not None and sa:
            ol, ot, ow, oh = shape_bounds(org)
            stx, sty, stw, sth = shape_bounds(sa[0])
            below = sty + sth
            sa_cx = stx + stw / 2
            # 组合形状构件：位于橙框内、Self-Attention 文本下方、横向靠近该列
            comp = [s for s in sh if s is not org
                    and ol <= center(s)[0] <= ol + ow and ot <= center(s)[1] <= ot + oh
                    and center(s)[1] >= below and abs(center(s)[0] - sa_cx) <= 1.3]
            if comp:
                dots = [s for s in comp if _is_dot(s) and
                        (color_is_orange(get_line_color(s)) or color_is_orange(get_fill_color(s)))]
                rects = [s for s in comp if looks_like_rect(s) and not _is_dot(s)
                         and shape_bounds(s)[2] < 2.5 and shape_bounds(s)[3] < 1.0]
                lines = [s for s in comp if _is_line_like2(s)]
                # ⑧⑨ 线条橙色、线宽0.5–1.25磅
                lines = [s for s in lines if color_is_orange(get_line_color(s))
                         and get_line_width_pt(s) is not None and in_range(get_line_width_pt(s), 0.5, 1.25)]
                vlines = [s for s in lines if is_vertical(s)]
                # ⑤ 三条虚线连接线；⑥ 三条虚线箭头
                dashed_conn = [s for s in lines if is_dash_line(s) and not has_arrow_head(s)]
                dashed_arrows = [s for s in lines if is_dash_line(s) and has_arrow_head(s)]
                # ⑦ 整体尺寸（组合全部构件包围盒）
                elems = dots + rects + lines
                if elems:
                    l = min(shape_bounds(s)[0] for s in elems)
                    r = max(shape_bounds(s)[0] + shape_bounds(s)[2] for s in elems)
                    t = min(shape_bounds(s)[1] for s in elems)
                    b = max(shape_bounds(s)[1] + shape_bounds(s)[3] for s in elems)
                    size_ok = in_range(r - l, 1.8, 2.4) and in_range(b - t, 1.7, 2.8)
                else:
                    size_ok = False
                att_ok = (len(dots) >= 1                      # ② 橙色节点圆点
                          and len(vlines) >= 1                # ③ 竖向线条
                          and len(rects) >= 1                 # ④ 底部小矩形
                          and len(dashed_conn) >= 3           # ⑤ 三条虚线连接线
                          and len(dashed_arrows) >= 3         # ⑥ 三条虚线箭头
                          and size_ok)                        # ⑦ 尺寸
        self.add(3, 'Self-Attention下方组合：橙色节点+竖线+底部矩形、3条虚线连线+3条虚线箭头、橙色0.5–1.25磅、整体1.8–2.4×1.7–2.8cm', att_ok)

        # 22. 第1页"Feed-Forward"文本下方组合形状。逐点检查：
        #    ① 位于 Feed-Forward 文本下方；② 6个橙色圆形节点；③ 10–16条橙色连接线；
        #    ④ 节点呈左三、中一、右二结构分布；⑤ 整体宽1.6–2.5cm、高1.6–2.8cm。
        def _is_line_like3(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))
        def _is_dot3(s):
            n = shape_name_or_prst(s)
            return 'ellipse' in n or 'oval' in n or 'circle' in n

        ff_ok = False
        ff = self.find_text('Feed-Forward')
        if org is not None and ff:
            ol, ot, ow, oh = shape_bounds(org)
            ftl, ftt, ftw, fth = shape_bounds(ff[0])
            below = ftt + fth
            ff_cx = ftl + ftw / 2
            comp = [s for s in sh if s is not org
                    and ol <= center(s)[0] <= ol + ow and ot <= center(s)[1] <= ot + oh
                    and center(s)[1] >= below and abs(center(s)[0] - ff_cx) <= 1.3]
            # ② 橙色圆形节点
            dots = [s for s in comp if _is_dot3(s)
                    and (color_is_orange(get_line_color(s)) or color_is_orange(get_fill_color(s)))]
            # ③ 橙色连接线：PPT里该组连线常使用浅橙棕色，按橙色系处理
            def _is_orange_family(c):
                if color_is_orange(c) or color_is_light_orange(c):
                    return True
                if c is None:
                    return False
                r, g, b = c[0], c[1], c[2]
                return r >= 210 and g >= 160 and b >= 130 and r >= g >= b and r - b >= 35
            lines = [s for s in comp if _is_line_like3(s) and _is_orange_family(get_line_color(s))]
            # ④ 左三、中一、右二：按节点中心 x 聚成的列及每列数量
            def _cluster(vals, tol=0.25):
                groups = []
                for v in sorted(vals):
                    if not groups or v - groups[-1][-1] > tol:
                        groups.append([v])
                    else:
                        groups[-1].append(v)
                return groups
            struct_ok = False
            if dots:
                cols = _cluster([center(s)[0] for s in dots])
                counts = [len(c) for c in cols]
                struct_ok = (counts == [3, 1, 2])
            # ⑤ 整体尺寸
            elems = dots + lines
            if elems:
                l = min(shape_bounds(s)[0] for s in elems)
                r = max(shape_bounds(s)[0] + shape_bounds(s)[2] for s in elems)
                t = min(shape_bounds(s)[1] for s in elems)
                b = max(shape_bounds(s)[1] + shape_bounds(s)[3] for s in elems)
                size_ok = in_range(r - l, 1.6, 2.5) and in_range(b - t, 1.6, 2.8)
            else:
                size_ok = False
            ff_ok = (len(dots) == 6                    # ② 6个节点
                     and 10 <= len(lines) <= 16        # ③ 10–16条连接线
                     and struct_ok                     # ④ 左三中一右二
                     and size_ok)                      # ⑤ 尺寸
        self.add(3, 'Feed-Forward下方组合：6个橙色圆节点(左三中一右二)+10–16条橙色连线、整体1.6–2.5×1.6–2.8cm', ff_ok)

        # 23. 第1页"Layer Normalization"文本下方组合图形/曲线图。逐点检查：
        #    ① 位于 Layer Normalization 文本下方；② 黑色水平轴；③ 黑色竖轴；
        #    ④ 橙色曲线；⑤ 曲线类似先上升后下降；⑥ 曲线线宽0.5–1.75磅；
        #    ⑦ 整体宽1.4–2cm、高1.5–2.2cm。
        def _is_line_like4(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))

        ln_ok = False
        lnorm = self.find_text('Layer Normalization')
        if org is not None and lnorm:
            ol, ot, ow, oh = shape_bounds(org)
            ltl, ltt, ltw, lth = shape_bounds(lnorm[0])
            below = ltt + lth
            ln_cx = ltl + ltw / 2
            comp = [s for s in sh if s is not org and _is_line_like4(s)
                    and ol <= center(s)[0] <= ol + ow and ot <= center(s)[1] <= ot + oh
                    and center(s)[1] >= below and abs(center(s)[0] - ln_cx) <= 1.3]
            # ② 黑色水平轴；③ 黑色竖轴
            haxis = [s for s in comp if color_is_black_or_dark(get_line_color(s)) and is_horizontal(s)]
            vaxis = [s for s in comp if color_is_black_or_dark(get_line_color(s)) and is_vertical(s)]
            # ④ 橙色曲线；⑥ 线宽0.5–1.75磅：曲线可能是描边线条，也可能是无描边/0线宽的 FREEFORM
            #    （PPT中该曲线用任意多边形实现），0线宽时改按线色或填充色是否橙色判断
            def _is_orange_curve(s):
                lw = get_line_width_pt(s)
                if lw is not None and lw > 0:
                    return color_is_orange(get_line_color(s)) and in_range(lw, 0.5, 1.75)
                return s.shape_type == MSO_SHAPE_TYPE.FREEFORM and (
                    color_is_orange(get_line_color(s)) or color_is_orange(get_fill_color(s)))
            curve = [s for s in comp if _is_orange_curve(s)]
            # ⑤ 先上升后下降：多段线条时按 x 排序检查中心 y 先减小(上升)后增大(下降)；
            #    若曲线是单个填充多边形（无法拆分线段），改为检查其自身形状的上边界
            #    是否呈"中间高、两端低"的先升后降轮廓（取形状顶点 y 坐标做同样的峰形判断）
            rise_fall = False
            if len(curve) >= 3:
                ys = [center(s)[1] for s in sorted(curve, key=lambda s: center(s)[0])]
                peak = ys.index(min(ys))  # 最高点(y最小)所在位置
                if 0 < peak < len(ys) - 1:
                    up = all(ys[i] >= ys[i + 1] for i in range(peak))          # 峰前上升
                    down = all(ys[i] <= ys[i + 1] for i in range(peak, len(ys) - 1))  # 峰后下降
                    rise_fall = up and down
            elif len(curve) >= 1:
                pts = freeform_points(curve[0])
                if pts and len(pts) >= 3:
                    xs_sorted = sorted(range(len(pts)), key=lambda i: pts[i][0])
                    ys = [pts[i][1] for i in xs_sorted]
                    peak = ys.index(min(ys))
                    if 0 < peak < len(ys) - 1:
                        # 曲线由多个采样点构成，允许局部小抖动（容差0.1cm），
                        # 只要求峰值前总体呈上升趋势、峰值后总体呈下降趋势
                        tol = 0.1
                        up = all(ys[i] >= ys[i + 1] - tol for i in range(peak))
                        down = all(ys[i] <= ys[i + 1] + tol for i in range(peak, len(ys) - 1))
                        rise_fall = up and down
            # ⑦ 整体尺寸
            elems = haxis + vaxis + curve
            if elems:
                l = min(shape_bounds(s)[0] for s in elems)
                r = max(shape_bounds(s)[0] + shape_bounds(s)[2] for s in elems)
                t = min(shape_bounds(s)[1] for s in elems)
                b = max(shape_bounds(s)[1] + shape_bounds(s)[3] for s in elems)
                size_ok = in_range(r - l, 1.4, 2.0) and in_range(b - t, 1.5, 2.2)
            else:
                size_ok = False
            ln_ok = (len(haxis) >= 1 and len(vaxis) >= 1 and len(curve) >= 1
                     and rise_fall and size_ok)
        self.add(3, 'Layer Normalization下方图形：黑色横轴+竖轴+橙色先升后降曲线(0.5–1.75磅)、整体1.4–2×1.5–2.2cm', ln_ok)

        # 24. 第1页橙色虚线矩形内部三个黑色箭头。逐点检查：
        #    ① 分别位于 Concatenate→Self-Attention、Self-Attention→Feed-Forward、
        #       Feed-Forward→Layer Normalization 四个模块矩形之间的三条间隙；
        #    ② 黑色；③ 水平；④ 箭头朝右；⑤ 线宽1–2磅；⑥ 单条长度0.4–1.1cm。
        def _is_line_like5(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))

        internal_ok = False
        if len(modules) == 4:
            # 四个模块矩形按 x 排序后，相邻两框之间的间隙区间
            mods_sorted = sorted(modules, key=lambda s: center(s)[0])
            gaps = []
            for i in range(3):
                lb = shape_bounds(mods_sorted[i])
                rb = shape_bounds(mods_sorted[i + 1])
                gap_l = lb[0] + lb[2]          # 左框右缘
                gap_r = rb[0]                  # 右框左缘
                gap_cy = (lb[1] + lb[3] / 2 + rb[1] + rb[3] / 2) / 2  # 两框中心 y 均值
                gaps.append((gap_l, gap_r, gap_cy))

            found = 0
            for gl, gr, gcy in gaps:
                for s in sh:
                    if not _is_line_like5(s):
                        continue
                    if not (color_is_black_or_dark(get_line_color(s)) and is_horizontal(s)):  # ②③
                        continue
                    if not horizontal_arrow_points_right(s):    # ④ 朝右
                        continue
                    lw = get_line_width_pt(s)
                    if lw is None or not in_range(lw, 1, 2):    # ⑤ 线宽1–2磅
                        continue
                    l, t, w, h = shape_bounds(s)
                    if not in_range(w, 0.4, 1.1):               # ⑥ 单条长度0.4–1.1cm
                        continue
                    cx, cy = center(s)
                    # ① 位于该间隙内、且竖向与模块中心线接近
                    if gl - 0.25 <= cx <= gr + 0.25 and abs(cy - gcy) <= 0.8:
                        found += 1
                        break
            internal_ok = (found == 3)
        self.add(3, '橙框内三个黑色水平右箭头：分居4个模块间的3条间隙、线宽1–2磅、单条长0.4–1.1cm', internal_ok)

        # 25. 第1页橙色虚线矩形右侧连接橙框与绿框的箭头。逐点检查：
        #    ① 位于距左27.8–31cm、距上4.2–5.6cm范围内；② 黑色；③ 水平；
        #    ④ 箭头朝右；⑤ 线宽1.2–2磅；⑥ 连接橙色虚线矩形与绿色虚线矩形。
        def _is_line_like6(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))
        grn = self.key.get('green_box')
        o2g_ok = False
        for s in sh:
            if not _is_line_like6(s):
                continue
            l, t, w, h = shape_bounds(s)
            # ① 范围
            if not (27.8 - 0.3 <= l and l + w <= 31 + 0.3 and 4.2 - 0.3 <= t and t + h <= 5.6 + 0.3):
                continue
            if not (color_is_black_or_dark(get_line_color(s)) and is_horizontal(s)):  # ②③
                continue
            if not horizontal_arrow_points_right(s):        # ④ 朝右
                continue
            lw = get_line_width_pt(s)
            if lw is None or not in_range(lw, 1.2, 2):       # ⑤ 线宽1.2–2磅
                continue
            # ⑥ 连接橙框与绿框：左端贴近橙框右缘、右端贴近绿框左缘
            connect_ok = True
            if org is not None:
                connect_ok = connect_ok and abs(l - (shape_bounds(org)[0] + shape_bounds(org)[2])) <= 0.6
            if grn is not None:
                connect_ok = connect_ok and abs((l + w) - shape_bounds(grn)[0]) <= 0.6
            if connect_ok:
                o2g_ok = True
                break
        self.add(3, '橙框右侧黑色水平右箭头(线宽1.2–2磅)，连接橙色虚线框与绿色虚线框(距左27.8–31/距上4.2–5.6cm)', o2g_ok)

        # 26. 第1页右上绿色矩形形状。逐点检查：① 位于距左29.2–36.5cm、距上1.6–11cm
        #    （已由 detect_key_objects 的 rrects 约束）；② 圆角矩形；③ 虚线；④ 绿色边线；
        #    ⑤ 宽4.5–5.5cm、高6.8–8.3cm（已由 rrects 约束）；⑥ 线宽0.75–2磅；
        #    ⑦ 内部无填充或浅绿色（以办公软件实际渲染为准：近白视为无填充）。
        green = self.key.get('green_box')
        green_ok = False
        if green is not None:
            lw = get_line_width_pt(green)
            green_ok = (is_rounded_rect(green)                          # 圆角矩形
                        and self.is_dashed_outline(green)               # 虚线
                        and color_is_green(get_line_color(green))       # 绿色边线
                        and lw is not None and in_range(lw, 0.75, 2)    # 线宽0.75–2磅
                        and (fill_is_visually_empty(green)              # 无填充（含近白视觉无填充）
                             or color_is_light_green(get_fill_color(green))))  # 或浅绿色
        strict_green_ok = green_ok
        self.add(5, '右上绿色虚线圆角矩形：位置尺寸达标、线宽0.75–2磅、无填充或浅绿色填充', green_ok)

        # 27. 第1页右上"4. Anomaly Scoring"文本。逐点检查：
        #    ① 位于绿色虚线矩形上方；② 文本仅1行；③ 字体 Arial 或 Calibri；
        #    ④ 字号12–18磅；⑤ 加粗；⑥ 颜色绿色。（细则未要求水平居中，故不加约束。）
        anom_title_ok = False
        if green is not None:
            gl, gt, gw, gh = shape_bounds(green)
            for s in self.find_text('4. Anomaly Scoring'):
                # ② 仅1行
                if rendered_line_count(s) != 1:
                    continue
                sl, st, sw, sh2 = shape_bounds(s)
                cx, cy = center(s)
                # ① 位于绿框上方：文本中心在框顶边以上，且横向与框重叠
                horiz_overlap = min(sl + sw, gl + gw) - max(sl, gl) > 0
                if not (cy <= gt + 0.1 and horiz_overlap):
                    continue
                # ③④⑤⑥ Arial/Calibri、12–18磅、加粗、绿色
                if text_font_ok(s, fonts=('arial', 'calibri'), size=(12, 18), bold=True,
                                color_pred=color_is_green):
                    anom_title_ok = True
                    break
        self.add(3, '绿框上方"4. Anomaly Scoring"：仅1行、Arial/Calibri 12–18磅加粗绿色', anom_title_ok)

        # 28. 第1页右上绿色虚线矩形内部偏上的矩形。逐点检查：
        #    ① 位于距左29.8–36cm、距上2.8–6cm范围内；② 圆角矩形；③ 浅绿色填充；
        #    ④ 边线绿色单实线；⑤ 线宽0.5–2磅；⑥ 内部文本"Reconstruction Module"；
        #    ⑦ 字号10–15磅；⑧ 颜色黑色；⑨ 文本分两行；⑩ 居中对齐。
        rec_ok = False
        for s in self.find_text('Reconstruction Module'):
            # ⑨ 两行；⑩ 居中；⑦ 字号10–15磅；⑧ 黑色
            if rendered_line_count(s) != 2:
                continue
            if not text_font_ok(s, size=(10, 15), color_pred=color_is_black_or_dark, center_align=True):
                continue
            # ②③④⑤ 承载文本的圆角矩形：浅绿填充、绿色单实线、线宽0.5–2磅
            def _rect_pred(r):
                lw = get_line_width_pt(r)
                return (color_is_light_green(get_fill_color(r))         # ③ 浅绿填充
                        and color_is_green(get_line_color(r))           # 绿色边线
                        and line_is_single_solid(r)                     # ④ 单实线
                        and not self.is_dashed_outline(r)               # 非虚线
                        and lw is not None and in_range(lw, 0.5, 2))    # ⑤ 线宽0.5–2磅
            rect = self.rect_covering_text(s, _rect_pred, margin=0.8)   # ② 圆角矩形
            # ① 矩形位于指定范围
            if rect is not None and box_in(rect, left=(29.8, 36), top=(2.8, 6)):
                rec_ok = True
                break
        self.add(3, '绿框内偏上浅绿圆角矩形(绿色单实线0.5–2磅)，内含两行10–15磅黑色居中"Reconstruction Module"', rec_ok)

        # 29. 第1页"Reconstruction Module"文本下方四个小方块。逐点检查：
        #    ① 位于 Reconstruction Module 文本下方；② 恰好4个；③ 横向排列且同一水平线；
        #    ④ 大小一致；⑤ 间距相等；⑥ 填充浅绿色或无填充；⑦ 边线绿色单实线；⑧ 线宽0.5–1磅。
        square_ok = False
        recs = self.find_text('Reconstruction Module')
        if green is not None and recs:
            gl, gt, gw, gh = shape_bounds(green)
            rl, rt, rw, rh = shape_bounds(recs[0])
            below = rt + rh
            rec_cx = rl + rw / 2
            squares = []
            for s in sh:
                if not looks_like_rect(s):
                    continue
                cx, cy = center(s)
                if not (gl <= cx <= gl + gw and below <= cy <= gt + gh):   # 绿框内、文本下方
                    continue
                bw2, bh2 = shape_bounds(s)[2], shape_bounds(s)[3]
                if not (bw2 < 1.2 and bh2 < 1.2):                          # 小方块
                    continue
                lw = get_line_width_pt(s)
                if not (color_is_green(get_line_color(s))                  # ⑦ 绿色边线
                        and line_is_single_solid(s)                        #    单实线
                        and lw is not None and in_range(lw, 0.5, 1)        # ⑧ 线宽0.5–1磅
                        and (color_is_light_green(get_fill_color(s))       # ⑥ 浅绿填充
                             or fill_is_visually_empty(s))):               #    或无填充
                    continue
                squares.append(s)
            if len(squares) == 4:                                          # ② 恰好4个
                squares.sort(key=lambda s: center(s)[0])
                cys = [center(s)[1] for s in squares]
                cxs = [center(s)[0] for s in squares]
                ws = [shape_bounds(s)[2] for s in squares]
                hs = [shape_bounds(s)[3] for s in squares]
                same_line = max(cys) - min(cys) < 0.15                     # ③ 同水平线
                size_same = (max(ws) - min(ws) <= 0.1) and (max(hs) - min(hs) <= 0.1)  # ④ 大小一致
                gaps = [cxs[i + 1] - cxs[i] for i in range(3)]
                equal_gap = (max(gaps) - min(gaps)) <= 0.12                # ⑤ 间距相等
                square_ok = same_line and size_same and equal_gap
        self.add(3, 'Reconstruction Module下方4个小方块：同水平、大小一致、等距、浅绿/无填充、绿色单实线0.5–1磅', square_ok)

        # 30. 第1页右上绿色虚线矩形内部偏下的矩形。逐点检查：
        #    ① 位于距左29.8–36cm、距上6.5–9.8cm范围内；② 圆角矩形；③ 浅绿色填充；
        #    ④ 边线绿色单实线；⑤ 线宽0.5–2磅；⑥ 内部文本"Anomaly Score"；
        #    ⑦ 文本仅1行；⑧ 字号10–15磅；⑨ 颜色黑色。
        score_ok = False
        for s in self.find_text('Anomaly Score'):
            # 排除标题"Anomaly Scoring"
            if text_contains(s, 'Anomaly Scoring'):
                continue
            # ⑦ 仅1行；⑧ 字号10–15磅；⑨ 黑色
            if rendered_line_count(s) != 1:
                continue
            if not text_font_ok(s, size=(10, 15), color_pred=color_is_black_or_dark):
                continue
            # ②③④⑤ 承载文本的圆角矩形：浅绿填充、绿色单实线、线宽0.5–2磅
            def _rect_pred(r):
                lw = get_line_width_pt(r)
                return (color_is_light_green(get_fill_color(r))         # ③ 浅绿填充
                        and color_is_green(get_line_color(r))           # 绿色边线
                        and line_is_single_solid(r)                     # ④ 单实线
                        and not self.is_dashed_outline(r)               # 非虚线
                        and lw is not None and in_range(lw, 0.5, 2))    # ⑤ 线宽0.5–2磅
            rect = self.rect_covering_text(s, _rect_pred, margin=0.8)   # ② 圆角矩形
            # ① 矩形位于指定范围
            if rect is not None and box_in(rect, left=(29.8, 36), top=(6.5, 9.8)):
                score_ok = True
                break
        self.add(3, '绿框内偏下浅绿圆角矩形(绿色单实线0.5–2磅)，内含一行10–15磅黑色"Anomaly Score"', score_ok)

        # 31. 第1页"Anomaly Score"下方曲线图标。逐点检查（细则每一个点都必须踩到，
        #    办公软件实际渲染为准）：
        #    ① 位于 Anomaly Score 文本下方；
        #    ② 黑色钟形/分布曲线：必须是"单一形状"或"整体为一张图片"，或者由多段
        #       小曲线拼接而成 —— 拼接时只允许一条大曲线（=一个连通分量），
        #       组成它的小段必须首尾相连而非分散；
        #    ③ 右侧红色高亮区域或红色曲线段；
        #    ④ 线宽0.5–2磅；
        #    ⑤ 图标整体宽2–3cm、高1–1.8cm。
        def _is_line_like7(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))
        def _is_red(c):
            return c is not None and c[0] > 150 and c[1] < 120 and c[2] < 120
        def _is_curve_shape(s):
            """单一形状承载的曲线：FREEFORM 路径含 arcTo 弧段、或路径含 quadBezTo/cubicBezTo。"""
            try:
                sp = s._element.spPr
                if sp is None:
                    return False
                for tag in ('a:arcTo', 'a:quadBezTo', 'a:cubicBezTo'):
                    if sp.find('.//' + qn(tag)) is not None:
                        return True
            except Exception:
                pass
            return False
        def _segment_endpoints(s):
            """按 xfrm + flipH/flipV 还原线段真实起终点，返回 [(x1,y1),(x2,y2)]。"""
            l, t, w, h = shape_bounds(s)
            xfrm = s._element.find('.//' + qn('a:xfrm'))
            flip_h = xfrm is not None and xfrm.get('flipH') == '1'
            flip_v = xfrm is not None and xfrm.get('flipV') == '1'
            if flip_h and flip_v:
                return [(l + w, t + h), (l, t)]
            if flip_h:
                return [(l + w, t), (l, t + h)]
            if flip_v:
                return [(l, t + h), (l + w, t)]
            return [(l, t), (l + w, t + h)]
        def _endpoints_connected(a, b, tol=0.05):
            """两段线是否首尾相连：任一端点距离 ≤ tol cm。"""
            ea = _segment_endpoints(a)
            eb = _segment_endpoints(b)
            for pa in ea:
                for pb in eb:
                    if abs(pa[0] - pb[0]) <= tol and abs(pa[1] - pb[1]) <= tol:
                        return True
            return False
        def _connected_components(shapes, tol=0.05):
            n = len(shapes)
            parent = list(range(n))
            def find(x):
                while parent[x] != x:
                    parent[x] = parent[parent[x]]
                    x = parent[x]
                return x
            for i in range(n):
                for j in range(i + 1, n):
                    if _endpoints_connected(shapes[i], shapes[j], tol):
                        a, b = find(i), find(j)
                        if a != b:
                            parent[a] = b
            return len(set(find(i) for i in range(n)))

        curve_ok = False
        scs = [s for s in self.find_text('Anomaly Score') if not text_contains(s, 'Anomaly Scoring')]
        if scs:
            stl, stt, stw, sth = shape_bounds(scs[0])
            below = stt + sth
            sc_cx = stl + stw / 2
            # 图标搜索窗口：紧贴 Anomaly Score 文本正下方
            window = [s for s in sh
                      if abs(center(s)[0] - sc_cx) <= 3.0
                      and center(s)[1] >= below and center(s)[1] <= below + 2.0]
            # ①④ 黑色/红色线条 且线宽0.5–2磅
            lines = [s for s in window if _is_line_like7(s)
                     and get_line_width_pt(s) is not None and in_range(get_line_width_pt(s), 0.5, 2)
                     and (color_is_black_or_dark(get_line_color(s)) or _is_red(get_line_color(s)))]
            # 排除水平基线（几乎 0 高，长度接近图标宽度）——它不是"钟形曲线"本身
            def _looks_like_baseline(s):
                l, t, w, h = shape_bounds(s)
                return h < 0.05 and w > 1.5
            black_lines = [s for s in lines
                           if color_is_black_or_dark(get_line_color(s)) and not _looks_like_baseline(s)]
            red_lines = [s for s in lines if _is_red(get_line_color(s))]

            # ② 黑色钟形曲线：单一形状（PICTURE/FREEFORM+曲线控制点），
            #   或多段拼接 = 一个连通分量（首尾相连）
            has_black_curve = False
            if black_lines:
                pics = [s for s in window
                        if getattr(s, 'shape_type', None) == MSO_SHAPE_TYPE.PICTURE]
                single_curve = [s for s in black_lines if _is_curve_shape(s)]
                if pics:
                    has_black_curve = True
                elif len(black_lines) == 1 and _is_curve_shape(black_lines[0]):
                    has_black_curve = True
                elif single_curve:
                    has_black_curve = True
                else:
                    # 拼接：所有黑色段必须首尾相连，形成"一条大曲线"
                    n_components = _connected_components(black_lines)
                    has_black_curve = (n_components == 1)

            has_red = len(red_lines) >= 1
            red_on_right = False
            if has_black_curve and has_red:
                black_cx = [center(s)[0] for s in black_lines]
                mid = (min(black_cx) + max(black_cx)) / 2
                red_on_right = (sum(center(s)[0] for s in red_lines) / len(red_lines)) >= mid
            # ⑤ 整体尺寸（黑+红曲线合并包围盒）
            icon_shapes = black_lines + red_lines
            if icon_shapes:
                cl = min(shape_bounds(s)[0] for s in icon_shapes)
                cr = max(shape_bounds(s)[0] + shape_bounds(s)[2] for s in icon_shapes)
                ct = min(shape_bounds(s)[1] for s in icon_shapes)
                cb = max(shape_bounds(s)[1] + shape_bounds(s)[3] for s in icon_shapes)
                size_ok = in_range(cr - cl, 2.0, 3.0) and in_range(cb - ct, 1.0, 1.8)
            else:
                size_ok = False
            curve_ok = has_black_curve and has_red and red_on_right and size_ok
        self.add(3, 'Anomaly Score下方曲线图标：单一/连通拼接的黑色分布曲线+右侧红色高亮/曲线段(0.5–2磅)、整体2–3×1–1.8cm', curve_ok)

        # 32. 第1页右上绿色虚线矩形内部、两个绿色实线矩形之间的箭头。逐点检查：
        #    ① 位于绿框内部两个绿色实线圆角矩形之间；② 黑色；③ 竖向；
        #    ④ 箭头朝下；⑤ 线宽0.5–2磅；⑥ 长度0.5–1cm。
        def _is_line_like8(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))
        green_arrow_ok = False
        if green is not None:
            gl, gt, gw, gh = shape_bounds(green)
            # 绿框内部的绿色实线圆角矩形（Reconstruction Module / Anomaly Score 两个承载框）
            green_rects = [s for s in sh if s is not green and looks_like_rect(s) and is_rounded_rect(s)
                           and color_is_green(get_line_color(s)) and line_is_single_solid(s)
                           and not self.is_dashed_outline(s)
                           and gl <= center(s)[0] <= gl + gw and gt <= center(s)[1] <= gt + gh]
            for s in sh:
                if not _is_line_like8(s):
                    continue
                if not (color_is_black_or_dark(get_line_color(s)) and is_vertical(s)):  # ②③
                    continue
                if not vertical_arrow_points_down(s):        # ④ 朝下
                    continue
                if not line_width_in_range(s, 0.5, 2):       # ⑤ 线宽0.5–2磅（未显式设线宽按默认0.75磅代入）
                    continue
                l, t, w, h = shape_bounds(s)
                if not in_range(h, 0.5, 1):                  # ⑥ 长度0.5–1cm
                    continue
                cx, cy = center(s)
                if not (gl <= cx <= gl + gw and gt <= cy <= gt + gh):  # 在绿框内
                    continue
                # ① 位于两个绿色实线矩形之间：上方有一个、下方有一个绿色矩形，且横向对齐
                above = any(center(r)[1] < t and abs(center(r)[0] - cx) <= 1.5 for r in green_rects)
                below_r = any(center(r)[1] > t + h and abs(center(r)[0] - cx) <= 1.5 for r in green_rects)
                if above and below_r:
                    green_arrow_ok = True
                    break
        self.add(3, '绿框内两个绿色实线矩形之间黑色竖向下箭头(线宽0.5–2磅、长0.5–1cm)', green_arrow_ok)

        # 33. 第1页右下粉色矩形形状。逐点检查：① 位于距左29.5–36.5cm、距上10.5–16.5cm
        #    （已由 detect_key_objects 的 rrects 约束）；② 圆角矩形；③ 虚线；④ 粉色边线；
        #    ⑤ 宽4.5–5.3cm、高4.5–5.5cm（已由 rrects 约束）；⑥ 线宽0.5–2磅；
        #    ⑦ 填充为无填充或粉色（以办公软件实际渲染为准：近白视为无填充）。
        pink = self.key.get('pink_box')
        pink_ok = False
        if pink is not None:
            lw = get_line_width_pt(pink)
            pink_ok = (is_rounded_rect(pink)                           # 圆角矩形
                       and self.is_dashed_outline(pink)                # 虚线
                       and color_is_pink(get_line_color(pink))         # 粉色边线
                       and lw is not None and in_range(lw, 0.5, 2)     # 线宽0.5–2磅
                       and (fill_is_visually_empty(pink)               # 无填充（含近白视觉无填充）
                            or color_is_pink(get_fill_color(pink))))   # 或粉色
        self.add(5, '右下粉色虚线圆角矩形：位置尺寸达标、线宽0.5–2磅、无填充或粉色填充', pink_ok)

        # 34. 第1页右下"5. Decision"文本。逐点检查：
        #    ① 位于粉色虚线矩形内部偏上；② 字体 Arial 或 Calibri；③ 字号12–16磅；
        #    ④ 加粗；⑤ 颜色粉色；⑥ 文本仅1行。
        decision_ok = False
        if pink is not None:
            pl, pt_, pw, ph = shape_bounds(pink)
            for s in self.find_text('5. Decision'):
                # ⑥ 仅1行
                if rendered_line_count(s) != 1:
                    continue
                cx, cy = center(s)
                # ① 位于粉框内部偏上（中心在框内且落在上半部）
                inside = pl <= cx <= pl + pw and pt_ <= cy <= pt_ + ph
                at_top = cy <= pt_ + ph / 2
                if not (inside and at_top):
                    continue
                # ②③④⑤ Arial/Calibri、12–16磅、加粗、粉色
                if text_font_ok(s, fonts=('arial', 'calibri'), size=(12, 16), bold=True,
                                color_pred=color_is_pink):
                    decision_ok = True
                    break
        self.add(3, '粉框内偏上"5. Decision"：仅1行、Arial/Calibri 12–16磅加粗粉色', decision_ok)

        # 35. 第1页粉色虚线矩形内部的一个矩形（Thresholding）。逐点检查：
        #    ① 位于距左31–35.5cm、距上12.1–14.4cm范围内；② 圆角矩形；
        #    ③ 白色或浅粉白色填充；④ 边线粉色单实线；⑤ 线宽0.5–1.75磅；
        #    ⑥ 内部文本"Thresholding"；⑦ 字号8–15磅；⑧ 颜色黑色；⑨ 加粗；⑩ 文本仅1行。
        th_ok = False
        for s in self.find_text('Thresholding'):
            # ⑩ 仅1行；⑦ 字号8–15磅；⑧ 黑色；⑨ 加粗
            if rendered_line_count(s) != 1:
                continue
            if not text_font_ok(s, size=(8, 15), bold=True, color_pred=color_is_black_or_dark):
                continue
            # ②③④⑤ 承载文本的圆角矩形：白/浅粉白填充、粉色单实线、线宽0.5–1.75磅
            def _rect_pred(r):
                lw = get_line_width_pt(r)
                return (color_is_white_or_near(get_fill_color(r))       # ③ 白/浅粉白填充
                        and color_is_pink(get_line_color(r))            # 粉色边线
                        and line_is_single_solid(r)                     # ④ 单实线
                        and not self.is_dashed_outline(r)               # 非虚线
                        and lw is not None and in_range(lw, 0.5, 1.75)) # ⑤ 线宽0.5–1.75磅
            rect = self.rect_covering_text(s, _rect_pred, margin=0.8)   # ② 圆角矩形
            # ① 矩形位于指定范围
            if rect is not None and box_in(rect, left=(31, 35.5), top=(12.1, 14.4)):
                th_ok = True
                break
        self.add(3, '粉框内白/浅粉白圆角矩形(粉色单实线0.5–1.75磅)，内含一行8–15磅加粗黑色"Thresholding"', th_ok)
        # 36. 第1页粉色虚线矩形内部的另一个矩形（Normal / Anomaly）。逐点检查：
        #    ① 位于距左31–35.5cm、距上14–15.8cm范围内；② 圆角矩形；
        #    ③ 白色或浅粉白色填充；④ 边线粉色单实线；⑤ 线宽0.5–1.75磅；
        #    ⑥ 内部文本"Normal / Anomaly"；⑦ 字号8–15磅；⑧ 颜色黑色；⑨ 加粗；⑩ 文本仅1行。
        #    "白色或浅粉白色填充"覆盖办公软件可能的所有呈现：真正白/近白、无填充(视觉白)、极浅粉。
        def _white_or_pinkish_light(shape):
            if fill_is_visually_empty(shape):        # 无填充/透明/近白，视觉上为白
                return True
            c = get_fill_color(shape)
            if c is None:
                return True
            if color_is_white_or_near(c):            # 白/近白
                return True
            # 浅粉白：高亮度、带极淡的粉红倾向（r 略高、整体很亮）
            r, g, b = c[0], c[1], c[2]
            return r >= 235 and g >= 220 and b >= 225 and r >= g

        na_ok = False
        for s in sh:
            # ⑥ 文本包含 Normal / Anomaly（兼容全角/半角斜杠、空格差异）
            txt = canonical_text(shape_text(s)).replace('／', '/')
            if 'normal/anomaly' not in txt:
                continue
            # ⑩ 仅1行；⑦ 字号8–15磅；⑧ 黑色；⑨ 加粗
            if rendered_line_count(s) != 1:
                continue
            if not text_font_ok(s, size=(8, 15), bold=True, color_pred=color_is_black_or_dark):
                continue
            # ②③④⑤ 承载文本的圆角矩形：白/浅粉白填充、粉色单实线、线宽0.5–1.75磅
            def _rect_pred(r):
                lw = get_line_width_pt(r)
                return (_white_or_pinkish_light(r)                     # ③ 白/浅粉白填充
                        and color_is_pink(get_line_color(r))           # 粉色边线
                        and line_is_single_solid(r)                    # ④ 单实线
                        and not self.is_dashed_outline(r)              # 非虚线
                        and lw is not None and in_range(lw, 0.5, 1.75))# ⑤ 线宽0.5–1.75磅
            rect = self.rect_covering_text(s, _rect_pred, margin=0.8)  # ② 圆角矩形
            # ① 矩形位于指定范围
            if rect is not None and box_in(rect, left=(31, 35.5), top=(14, 15.8)):
                na_ok = True
                break
        self.add(3, '粉框内白/浅粉白圆角矩形(粉色单实线0.5–1.75磅)，内含一行8–15磅加粗黑色"Normal / Anomaly"', na_ok)

        # 37. 第1页粉色虚线矩形内部、两个粉色实线矩形之间的箭头。逐点检查：
        #    ① 位于粉框内部两个粉色实线圆角矩形之间；② 黑色；③ 竖向；
        #    ④ 箭头朝下；⑤ 线宽0.5–2磅；⑥ 长度0.2–0.6cm。
        def _is_line_like9(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))
        pink_arrow_ok = False
        if pink is not None:
            pl, pt_, pw, ph = shape_bounds(pink)
            # 粉框内部的粉色实线圆角矩形（Thresholding / Normal-Anomaly 两个承载框）
            pink_rects = [s for s in sh if s is not pink and looks_like_rect(s) and is_rounded_rect(s)
                          and color_is_pink(get_line_color(s)) and line_is_single_solid(s)
                          and not self.is_dashed_outline(s)
                          and pl <= center(s)[0] <= pl + pw and pt_ <= center(s)[1] <= pt_ + ph]
            for s in sh:
                if not _is_line_like9(s):
                    continue
                if not (color_is_black_or_dark(get_line_color(s)) and is_vertical(s)):  # ②③
                    continue
                if not vertical_arrow_points_down(s):        # ④ 朝下
                    continue
                lw = get_line_width_pt(s)
                if lw is None or not in_range(lw, 0.5, 2):   # ⑤ 线宽0.5–2磅
                    continue
                l, t, w, h = shape_bounds(s)
                if not in_range(h, 0.2, 0.6):                # ⑥ 长度0.2–0.6cm
                    continue
                cx, cy = center(s)
                if not (pl <= cx <= pl + pw and pt_ <= cy <= pt_ + ph):  # 在粉框内
                    continue
                # ① 位于两个粉色实线矩形之间：上方有一个、下方有一个粉色矩形，横向对齐
                above = any(center(r)[1] < t and abs(center(r)[0] - cx) <= 1.5 for r in pink_rects)
                below_r = any(center(r)[1] > t + h and abs(center(r)[0] - cx) <= 1.5 for r in pink_rects)
                if above and below_r:
                    pink_arrow_ok = True
                    break
        self.add(3, '粉框内两个粉色实线矩形之间黑色竖向下箭头(线宽0.5–2磅、长0.2–0.6cm)', pink_arrow_ok)
        # 38. 第1页绿色虚线矩形与粉色虚线矩形之间的竖向箭头。逐点检查：
        #    ① 位于距左32–34cm、距上9.7–11.6cm范围内；② 黑色；③ 竖向；④ 箭头朝下；
        #    ⑤ 线宽1.2–2磅；⑥ 长度1.2–1.6cm；⑦ 连接绿色虚线矩形与粉色虚线矩形。
        def _is_line_like10(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))
        g2p_ok = False
        for s in sh:
            if not _is_line_like10(s):
                continue
            l, t, w, h = shape_bounds(s)
            # ① 范围
            if not (32 - 0.3 <= l and l + w <= 34 + 0.3 and 9.7 - 0.3 <= t and t + h <= 11.6 + 0.3):
                continue
            if not (color_is_black_or_dark(get_line_color(s)) and is_vertical(s)):  # ②③
                continue
            if not vertical_arrow_points_down(s):        # ④ 朝下
                continue
            lw = get_line_width_pt(s)
            if lw is None or not in_range(lw, 1.2, 2):   # ⑤ 线宽1.2–2磅
                continue
            if not in_range(h, 1.2, 1.6):                # ⑥ 长度1.2–1.6cm
                continue
            # ⑦ 连接绿框与粉框：上端贴近绿框底缘、下端贴近粉框顶缘
            connect_ok = True
            if green is not None:
                connect_ok = connect_ok and abs(t - (shape_bounds(green)[1] + shape_bounds(green)[3])) <= 0.6
            if pink is not None:
                connect_ok = connect_ok and abs((t + h) - shape_bounds(pink)[1]) <= 0.6
            if connect_ok:
                g2p_ok = True
                break
        self.add(3, '绿框与粉框之间黑色竖向下箭头(线宽1.2–2磅、长1.2–1.6cm)，连接两框(距左32–34/距上9.7–11.6cm)', g2p_ok)

        # 39. 第1页下方蓝色矩形形状。逐点检查：① 位于距左15.0–25.3cm、距上10.0–16cm
        #    （已由 detect_key_objects 的 rrects 约束）；② 圆角矩形；③ 虚线；④ 蓝色边线；
        #    ⑤ 宽8–9.5cm、高4.4–5cm（已由 rrects 约束）；⑥ 线宽0.5–2磅；
        #    ⑦ 填充为无填充或浅蓝色（以办公软件实际渲染为准：近白视为无填充）。
        blue = self.key.get('blue_box')
        blue_ok = False
        if blue is not None:
            lw = get_line_width_pt(blue)
            blue_ok = (is_rounded_rect(blue)                          # 圆角矩形
                       and self.is_dashed_outline(blue)               # 虚线
                       and color_is_blue(get_line_color(blue))        # 蓝色边线
                       and lw is not None and in_range(lw, 0.5, 2)    # 线宽0.5–2磅
                       and (fill_is_visually_empty(blue)              # 无填充（含近白视觉无填充）
                            or color_is_light_blue(get_fill_color(blue))))  # 或浅蓝色
        self.add(5, '下方蓝色虚线圆角矩形：位置尺寸达标、线宽0.5–2磅、无填充或浅蓝色填充', blue_ok)

        # 40. 第1页蓝色虚线矩形内部"6. Training Objective"文本。逐点检查：
        #    ① 位于蓝色矩形顶部；② 水平居中；③ 文本仅1行；④ 字体 Arial 或 Calibri；
        #    ⑤ 字号12–18磅；⑥ 加粗；⑦ 颜色蓝色。
        def _training_objective_single_line(shape):
            """针对标题文本框严格判断视觉上是否只占一行。"""
            if rendered_line_count(shape) != 1:
                return False
            try:
                body = shape._element.find('.//' + qn('a:bodyPr'))
                wrap = 'square'
                l_ins = r_ins = 91440
                font_scale = 1.0
                if body is not None:
                    wrap = body.get('wrap', 'square') or 'square'
                    l_ins = int(body.get('lIns', str(l_ins)) or l_ins)
                    r_ins = int(body.get('rIns', str(r_ins)) or r_ins)
                    norm = body.find('.//' + qn('a:normAutofit'))
                    if norm is not None and norm.get('fontScale'):
                        font_scale = int(norm.get('fontScale')) / 100000.0
                paragraphs = [p for p in shape.text_frame.paragraphs
                              if ''.join(r.text for r in p.runs).strip()]
                if len(paragraphs) != 1:
                    return False
                if wrap == 'none':
                    return True
                avail_cm = emu2cm(shape.width - l_ins - r_ins)
                if avail_cm <= 0:
                    return False
                runs = _runs_of_paragraph(paragraphs[0])
                if font_scale != 1.0:
                    runs = [(t, n, b, sz * font_scale) for (t, n, b, sz) in runs]
                full_w = sum(_measure_text_cm(t, n, b, sz) for (t, n, b, sz) in runs)
                return full_w <= avail_cm
            except Exception:
                return False

        train_ok = False
        if blue is not None:
            bl, bt, bw, bh = shape_bounds(blue)
            for s in self.find_text('6. Training Objective'):
                # ③ 仅1行：显式换行和文本框自动换行都不能出现
                if not _training_objective_single_line(s):
                    continue
                cx, cy = center(s)
                # ① 位于蓝框顶部：中心在框内且落在上半部
                inside = bl <= cx <= bl + bw and bt <= cy <= bt + bh
                at_top = cy <= bt + bh / 2
                # ② 水平居中：文本中心 x 接近框中线（±0.8cm），或段落设为居中
                centered_pos = abs(cx - (bl + bw / 2)) <= 0.8
                if not (inside and at_top):
                    continue
                # ②(续)④⑤⑥⑦ 居中(段落对齐或位置居中)、Arial/Calibri、12–18磅、加粗、蓝色
                font_ok = text_font_ok(s, fonts=('arial', 'calibri'), size=(12, 18), bold=True,
                                       color_pred=color_is_blue)
                align_center = text_font_ok(s, fonts=('arial', 'calibri'), size=(12, 18), bold=True,
                                            color_pred=color_is_blue, center_align=True)
                if font_ok and (align_center or centered_pos):
                    train_ok = True
                    break
        self.add(3, '蓝框顶部居中"6. Training Objective"：仅1行、Arial/Calibri 12–18磅加粗蓝色', train_ok)

        # 41. 第1页蓝色虚线矩形形状内部出现一个矩形：位于"6. Training Objective"文本左下，
        #     距左16–20cm、距上11.5–14.5cm范围内。逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 位于蓝色虚线矩形（blue）内部；② 位于"6. Training Objective"文本左下（在其下方、
        #     不越过其右缘）；③ 位置范围：距左16–20cm、距上11.5–14.5cm；④ 形状为圆角矩形；
        #     ⑤ 浅蓝色填充；⑥ 边线为蓝色单实线；⑦ 线宽0.5–1.5磅；⑧ 其内部出现"Reconstruction Loss"
        #     和"Lrec"文本；⑨ 两文本分别位于独立的两行（各自独立一行、上下堆叠）；
        #     ⑩ 字号8–18磅；⑪ 颜色为黑色；⑫ "Lrec"中 rec 字体小于 L。
        def _blue_rect_pred(r):
            lw = get_line_width_pt(r)
            return (is_rounded_rect(r)                                      # ④ 圆角矩形
                    and color_is_light_blue(get_fill_color(r))              # ⑤ 浅蓝色填充
                    and color_is_blue(get_line_color(r))                    # ⑥ 蓝色边线
                    and line_is_single_solid(r)                             # ⑥ 单实线
                    and lw is not None and in_range(lw, 0.5, 1.5))          # ⑦ 线宽0.5–1.5磅

        rl_ok = False
        if blue is not None:
            bbl, bbt, bbw, bbh = shape_bounds(blue)
            to_refs = self.find_text('6. Training Objective')
            to_ref = to_refs[0] if to_refs else None
            for main in self.find_text('Reconstruction Loss'):
                # ⑨ 主文本独立一行；⑩⑪ 字号8–18磅、黑色（不校验字体/加粗/居中——细则未要求）
                if rendered_line_count(main) != 1:
                    continue
                if not text_font_ok(main, fonts=None, size=(8, 18),
                                    color_pred=color_is_black_or_dark):
                    continue
                rect = self.rect_covering_text(main, _blue_rect_pred, margin=0.8)
                if rect is None:
                    continue
                rl, rt, rw, rh = shape_bounds(rect)
                rcx, rcy = rl + rw / 2, rt + rh / 2
                # ① 承载矩形位于蓝色虚线矩形内部
                if not (bbl <= rcx <= bbl + bbw and bbt <= rcy <= bbt + bbh):
                    continue
                # ③ 位置范围
                if not box_in(rect, left=(16, 20), top=(11.5, 14.5)):
                    continue
                # ② 位于"6. Training Objective"文本左下（在其下方，且不越过其右缘）
                if to_ref is not None:
                    tl, tt, tw, th = shape_bounds(to_ref)
                    if not (rcy >= tt + th and rcx <= tl + tw):
                        continue
                # ⑧ 内部还需出现"Lrec"文本；⑨ 与主文本独立两行；⑩⑪ 8–18磅黑色；⑫ rec 小于 L
                mcy = center(main)[1]
                found_lrec = False
                for s in sh:
                    if s is main:
                        continue
                    cn = canonical_text(shape_text(s)).replace('_', '')
                    if 'lrec' not in cn:
                        continue
                    scx, scy = center(s)
                    # 位于同一承载矩形内
                    if not (rl <= scx <= rl + rw and rt <= scy <= rt + rh):
                        continue
                    # ⑨ 与主文本上下堆叠（不同行）；且各自独立一行
                    if abs(scy - mcy) < 0.05:
                        continue
                    if rendered_line_count(s) != 1:
                        continue
                    # ⑩⑪ 字号8–18磅、黑色
                    if not text_font_ok(s, fonts=None, size=(8, 18),
                                        color_pred=color_is_black_or_dark):
                        continue
                    # ⑫ "Lrec"中 rec 字体小于 L（办公软件实际渲染更小：不同run更小/下标/上标）
                    if not has_smaller_suffix_text(s, 'L', 'rec'):
                        continue
                    found_lrec = True
                    break
                if found_lrec:
                    rl_ok = True
                    break
        self.add(3, '蓝框内左下浅蓝圆角矩形(蓝色单实线0.5–1.5磅)，位于"6. Training Objective"左下(距左16–20/距上11.5–14.5cm)：两行8–18磅黑色"Reconstruction Loss"/"Lrec"，rec小于L', rl_ok)
        # 42. 第1页蓝色虚线矩形形状内部出现一个矩形：位于"6. Training Objective"文本右下，
        #     距左20–24.6cm、距上11.5–14.5cm范围内。逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 位于蓝色虚线矩形（blue）内部；② 位于"6. Training Objective"文本右下（在其下方、
        #     不越过其左缘，即整体位于该文本右侧下方）；③ 位置范围：距左20–24.6cm、距上11.5–14.5cm；
        #     ④ 形状为圆角矩形；⑤ 浅蓝色填充；⑥ 边线为蓝色单实线；⑦ 线宽0.5–1.5磅；
        #     ⑧ 其内部出现"Regularization Loss"和"Lreg"文本；⑨ 两文本分别位于独立的两行；
        #     ⑩ 字号8–18磅；⑪ 颜色为黑色；⑫ "Lreg"中 reg 字体小于 L。
        def _blue_rect_pred_reg(r):
            lw = get_line_width_pt(r)
            return (is_rounded_rect(r)                                      # ④ 圆角矩形
                    and color_is_light_blue(get_fill_color(r))              # ⑤ 浅蓝色填充
                    and color_is_blue(get_line_color(r))                    # ⑥ 蓝色边线
                    and line_is_single_solid(r)                             # ⑥ 单实线
                    and lw is not None and in_range(lw, 0.5, 1.5))          # ⑦ 线宽0.5–1.5磅

        reg_ok = False
        if blue is not None:
            bbl, bbt, bbw, bbh = shape_bounds(blue)
            to_refs = self.find_text('6. Training Objective')
            to_ref = to_refs[0] if to_refs else None
            for main in self.find_text('Regularization Loss'):
                # ⑨ 主文本独立一行；⑩⑪ 字号8–18磅、黑色（不校验字体/加粗/居中——细则未要求）
                if rendered_line_count(main) != 1:
                    continue
                if not text_font_ok(main, fonts=None, size=(8, 18),
                                    color_pred=color_is_black_or_dark):
                    continue
                rect = self.rect_covering_text(main, _blue_rect_pred_reg, margin=0.8)
                if rect is None:
                    continue
                rl, rt, rw, rh = shape_bounds(rect)
                rcx, rcy = rl + rw / 2, rt + rh / 2
                # ① 承载矩形位于蓝色虚线矩形内部
                if not (bbl <= rcx <= bbl + bbw and bbt <= rcy <= bbt + bbh):
                    continue
                # ③ 位置范围
                if not box_in(rect, left=(20, 24.6), top=(11.5, 14.5)):
                    continue
                # ② 位于"6. Training Objective"文本右下（在其下方，且不越过其左缘）
                if to_ref is not None:
                    tl, tt, tw, th = shape_bounds(to_ref)
                    if not (rcy >= tt + th and rcx >= tl):
                        continue
                # ⑧ 内部还需出现"Lreg"文本；⑨ 与主文本独立两行；⑩⑪ 8–18磅黑色；⑫ reg 小于 L
                mcy = center(main)[1]
                found_lreg = False
                for s in sh:
                    if s is main:
                        continue
                    cn = canonical_text(shape_text(s)).replace('_', '')
                    if 'lreg' not in cn:
                        continue
                    # 排除 Total Loss 公式那种"L = Lrec + λLreg"的长文本，公式在同一行里同时含 lrec 与 lreg
                    if 'lrec' in cn:
                        continue
                    scx, scy = center(s)
                    # 位于同一承载矩形内
                    if not (rl <= scx <= rl + rw and rt <= scy <= rt + rh):
                        continue
                    # ⑨ 与主文本上下堆叠（不同行）；且各自独立一行
                    if abs(scy - mcy) < 0.05:
                        continue
                    if rendered_line_count(s) != 1:
                        continue
                    # ⑩⑪ 字号8–18磅、黑色
                    if not text_font_ok(s, fonts=None, size=(8, 18),
                                        color_pred=color_is_black_or_dark):
                        continue
                    # ⑫ "Lreg"中 reg 字体小于 L（办公软件实际渲染更小：不同run更小/下标/上标）
                    if not has_smaller_suffix_text(s, 'L', 'reg'):
                        continue
                    found_lreg = True
                    break
                if found_lreg:
                    reg_ok = True
                    break
        self.add(3, '蓝框内右下浅蓝圆角矩形(蓝色单实线0.5–1.5磅)，位于"6. Training Objective"右下(距左20–24.6/距上11.5–14.5cm)：两行8–18磅黑色"Regularization Loss"/"Lreg"，reg小于L', reg_ok)

        # 43. 第1页蓝色虚线矩形形状内部出现"+"文本：位于两个蓝色实线矩形之间，
        #     距左19–21cm、距上12.2–13.6cm范围内，字体为 Arial 或 Calibri，字号14–18磅，颜色黑色。
        #     逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 位于蓝色虚线矩形内部；② 文本为"+"；③ 位于两个蓝色实线矩形之间；
        #     ④ 位置范围：距左19–21cm、距上12.2–13.6cm；
        #     ⑤ 字体 Arial 或 Calibri；⑥ 字号14–18磅；⑦ 颜色黑色。
        plus_ok = False
        if blue is not None:
            bbl, bbt, bbw, bbh = shape_bounds(blue)
            # 蓝色虚线框内部的两个蓝色实线圆角矩形（承载 Reconstruction/Regularization Loss）
            blue_solid_rects = [s for s in sh if s is not blue
                                and looks_like_rect(s) and is_rounded_rect(s)
                                and color_is_blue(get_line_color(s))
                                and line_is_single_solid(s)
                                and not self.is_dashed_outline(s)
                                and bbl <= center(s)[0] <= bbl + bbw
                                and bbt <= center(s)[1] <= bbt + bbh]
            for s in sh:
                # ② 文本为"+"
                if shape_text(s).strip() != '+':
                    continue
                pcx, pcy = center(s)
                # ① 位于蓝色虚线矩形内部
                if not (bbl <= pcx <= bbl + bbw and bbt <= pcy <= bbt + bbh):
                    continue
                # ④ 距左19–21cm、距上12.2–13.6cm（按办公软件里形状定位=左上角坐标）
                if not box_in(s, left=(19, 21), top=(12.2, 13.6)):
                    continue
                # ③ 位于两个蓝色实线矩形之间：至少存在一个在其左侧、一个在其右侧
                left_side = any(center(r)[0] < pcx for r in blue_solid_rects)
                right_side = any(center(r)[0] > pcx for r in blue_solid_rects)
                if not (len(blue_solid_rects) >= 2 and left_side and right_side):
                    continue
                # ⑤⑥⑦ 字体 Arial/Calibri、字号14–18磅、颜色黑色
                if not text_font_ok(s, fonts=('arial', 'calibri'),
                                    size=(14, 18), color_pred=color_is_black_or_dark):
                    continue
                plus_ok = True
                break
        self.add(3, '蓝框内两个蓝色实线矩形之间"+"文本：Arial/Calibri 14–18磅黑色(距左19–21/距上12.2–13.6cm)', plus_ok)

        # 44. 第1页蓝色虚线矩形形状内部出现蓝色直线：位于蓝色虚线矩形下部，长度7.3–7.8cm之间。
        #     逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 位于蓝色虚线矩形内部；② 直线；③ 蓝色；④ 位于蓝框下部；⑤ 长度7.3–7.8cm。
        def _is_line_like_blue(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))

        blue_line_ok = False
        if blue is not None:
            bbl, bbt, bbw, bbh = shape_bounds(blue)
            for s in sh:
                # ② 直线（线段/连接符/直线形状）
                if not _is_line_like_blue(s):
                    continue
                # ③ 蓝色
                if not color_is_blue(get_line_color(s)):
                    continue
                l, t, w, h = shape_bounds(s)
                cx, cy = l + w / 2, t + h / 2
                # ① 位于蓝色虚线矩形内部（线段中心落在蓝框内）
                if not (bbl <= cx <= bbl + bbw and bbt <= cy <= bbt + bbh):
                    continue
                # ④ 位于蓝框下部（线段中心在框下半区）
                if not (cy >= bbt + bbh / 2):
                    continue
                # ⑤ 长度7.3–7.8cm（按办公软件里可见的线段长度：水平/竖直取较长边，其它取对角）
                length = max(w, h) if (w < 0.05 or h < 0.05) else math.hypot(w, h)
                if not in_range(length, 7.3, 7.8):
                    continue
                blue_line_ok = True
                break
        self.add(1, '蓝框内部下部蓝色直线(长度7.3–7.8cm)', blue_line_ok)

        # 45. 第1页蓝色虚线矩形形状内部出现"Total Loss: L = Lrec + λLreg"文本或相近公式：
        #     位于蓝色虚线矩形内部底部居中，距左17–23.3cm、距上10.7–15.9cm范围内，
        #     文本或公式内容整体居于一行排列；字体 Arial 或 Calibri，字号9–18磅，颜色黑色。
        #     逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 位于蓝色虚线矩形内部；② 位于蓝框内部底部；③ 水平居中；
        #     ④ 距左17–23.3cm、距上10.7–15.9cm；⑤ 文本包含"Total Loss"及"Lrec""Lreg"等相近公式内容；
        #     ⑥ 整体位于一行；⑦ 字体 Arial 或 Calibri；⑧ 字号9–18磅；⑨ 颜色黑色。
        formula_ok = False
        if blue is not None:
            bbl, bbt, bbw, bbh = shape_bounds(blue)
            for s in sh:
                # ⑤ 文本包含 Total Loss + Lrec + Lreg（相近公式，兼容下划线/空格差异）
                txt = canonical_text(shape_text(s))
                if not (('totalloss' in txt and 'lrec' in txt and 'lreg' in txt)
                        or ('l=lrec' in txt and 'lreg' in txt)):
                    continue
                fl, ft, fw, fh = shape_bounds(s)
                fcx, fcy = fl + fw / 2, ft + fh / 2
                # ① 位于蓝色虚线矩形内部（文本中心落在蓝框内）
                if not (bbl <= fcx <= bbl + bbw and bbt <= fcy <= bbt + bbh):
                    continue
                # ② 位于蓝框内部底部（中心在下半区）
                if not (fcy >= bbt + bbh / 2):
                    continue
                # ③ 水平居中：文本中心 x 接近蓝框中线（±0.8cm）或段落设为居中
                center_pos = abs(fcx - (bbl + bbw / 2)) <= 0.8
                center_align_ok = False
                try:
                    aligns = [p.alignment for p in s.text_frame.paragraphs if p.text.strip()]
                    center_align_ok = any(
                        (a is not None) and (str(a).upper().endswith('CENTER') or a == 2)
                        for a in aligns
                    )
                except Exception:
                    pass
                if not (center_pos or center_align_ok):
                    continue
                # ④ 距左17–23.3cm、距上10.7–15.9cm
                if not box_in(s, left=(17, 23.3), top=(10.7, 15.9)):
                    continue
                # ⑥ 整体位于一行
                if rendered_line_count(s) != 1:
                    continue
                # ⑦⑧⑨ 字体 Arial/Calibri、字号9–18磅、颜色黑色
                if not text_font_ok(s, fonts=('arial', 'calibri'),
                                    size=(9, 18), color_pred=color_is_black_or_dark):
                    continue
                formula_ok = True
                break
        self.add(3, '蓝框内部底部居中"Total Loss: L = Lrec + λLreg"或相近公式：1行、Arial/Calibri 9–18磅黑色(距左17–23.3/距上10.7–15.9cm)', formula_ok)

        # 46-47. 到蓝色框的箭头：必须是单一形状承载的圆角折线，不接受多段直线/短线段拼接
        # 46. 第1页橙色虚线矩形到蓝色虚线矩形之间出现一个箭头：
        #     位于距左19.5–22.8cm、距上8.5–11.2cm范围内，为黑色圆角折线箭头，有两个拐点；
        #     线宽0.5–2磅，箭头方向为先向下再向左最后向下。
        #     逐点检查（细则每一个点都必须踩到，办公软件实际渲染为准；同一条箭头必须同时满足所有条件）：
        #     ① 位于距左19.5–22.8cm、距上8.5–11.2cm范围内；
        #     ② 黑色；
        #     ③ 线宽0.5–2磅；
        #     ④ 末端带箭头；
        #     ⑤ 圆角折线（bentConnector3/curvedConnector3、FREEFORM+arcTo、或折线+<a:round/>）；
        #     ⑥ 两个拐点（bentConnector3/curvedConnector3 天然=2 拐点；bentConnector2/4 排除）；
        #     ⑦ 方向：先向下→再向左→最后向下（上端起点、下端末端；箭头在下方左侧竖段）。
        orange_to_blue_arrow = False
        lx0, lx1, ty0, ty1 = 19.5, 22.8, 8.5, 11.2
        tol_r = 0.4
        for s in self.shapes:
            if not (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s)
                    or 'connector' in shape_name_or_prst(s)):
                continue
            if not line_color_is_black_or_dark(s):                # ② 黑色（含主题色/继承默认）
                continue
            if not line_width_in_range(s, 0.5, 2):                # ③ 线宽0.5–2磅（含继承主题/默认）
                continue
            l, t, w, h = shape_bounds(s)
            if not (lx0 - tol_r <= l and l + w <= lx1 + tol_r
                    and ty0 - tol_r <= t and t + h <= ty1 + tol_r):  # ①
                continue
            prst = ''
            try:
                g = s._element.spPr.find(qn('a:prstGeom'))
                if g is not None:
                    prst = g.get('prst', '') or ''
            except Exception:
                pass
            # ⑤⑥ 只允许"单一形状承载的圆角折线且天然两个拐点"：
            #     - bentConnector3 / curvedConnector3（天然=2拐点），或
            #     - FREEFORM 路径含 a:arcTo 弧段的圆角折线，或
            #     - 显式 <a:round/> 连接样式的折线形状。
            #     直线（line / straightConnector1）不具备"两个拐点+圆角折线"，一律排除。
            is_rounded = False
            if prst in ('bentConnector3', 'curvedConnector3'):
                if w >= 0.4 and h >= 0.4:                         # 两个拐点：bbox 同时跨水平+竖直
                    is_rounded = True
            elif (prst.startswith('bentConnector')
                  or prst.startswith('curvedConnector')):
                pass
            else:
                try:
                    if s.shape_type == MSO_SHAPE_TYPE.FREEFORM \
                            and s._element.spPr.find('.//' + qn('a:arcTo')) is not None \
                            and w >= 0.4 and h >= 0.4:
                        is_rounded = True
                except Exception:
                    pass
                if not is_rounded:
                    try:
                        ln = s._element.spPr.find(qn('a:ln'))
                        if ln is not None and ln.find(qn('a:round')) is not None \
                                and w >= 0.4 and h >= 0.4:
                            is_rounded = True
                    except Exception:
                        pass
            if not is_rounded:
                continue
            if not has_arrow_head(s):                             # ④
                continue
            # ⑦ 方向"下→左→下"：起点在右上、末端在左下 —— 从 bbox+flip 反演：
            #   未 flip 的 bentConnector3 起点(左上)→终点(右下)。
            #   "右上→左下" 要求 x 反向 → flipH=1；y 保持 → flipV=0。
            xfrm = s._element.find('.//' + qn('a:xfrm'))
            flip_h = xfrm is not None and xfrm.get('flipH') == '1'
            flip_v = xfrm is not None and xfrm.get('flipV') == '1'
            if not flip_h:
                continue
            if flip_v:
                continue
            orange_to_blue_arrow = True
            break
        self.add(3, '橙框到蓝框黑色圆角折线箭头：先下→左→下、两拐点、线宽0.5–2磅(距左19.5–22.8/距上8.5–11.2cm)', orange_to_blue_arrow)

        # 47. 第1页紫色虚线矩形到蓝色虚线矩形之间出现一个箭头：
        #     位于距左10.6–19.9cm、距上7.5–12.8cm范围内，为蓝色虚线圆角折线箭头，
        #     线宽0.5–2磅，箭头方向为先向左再向上。
        #     逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 位于距左10.6–19.9cm、距上7.5–12.8cm范围内；② 蓝色；③ 虚线；
        #     ④ 折线箭头（末端带箭头）；⑤ 线宽0.5–2磅；⑥ 方向：先向左→再向上；⑦ 圆角。
        def _one_turn_arrow_in_region(region, dirs, color_pred=None, width=None,
                                      dashed_only=True, tol=0.4):
            """检测一条只含一个拐点的两段折线箭头（首段+末段带箭头）。
            办公软件中虚线折线通常由多段共线短段拼成，先按方向把共线相邻短段合并成整段，
            再枚举两段首尾相连 + 方向序列 + 末段带箭头。
            细则要求"圆角折线"：要么由 bent/curved connector / FREEFORM+arcTo 单一形状承载，
            要么参与拼接的每一段折线自身携带 <a:round/> 连接样式；普通直线/straightConnector1
            拼出的直角折线不算圆角，一律排除。"""
            def _prst(s):
                try:
                    g = s._element.spPr.find(qn('a:prstGeom'))
                    if g is not None:
                        return g.get('prst', '') or ''
                except Exception:
                    pass
                return ''

            def is_rounded_seg(s):
                # bent/curved connector 单一形状承载圆角折线
                p = _prst(s)
                if p.startswith('bentConnector') or p.startswith('curvedConnector'):
                    return True
                # FREEFORM 路径含 arcTo 弧段
                try:
                    if s.shape_type == MSO_SHAPE_TYPE.FREEFORM \
                            and s._element.spPr.find('.//' + qn('a:arcTo')) is not None:
                        return True
                except Exception:
                    pass
                # 折线段显式 <a:round/> 连接样式
                try:
                    ln = s._element.spPr.find(qn('a:ln'))
                    if ln is not None and ln.find(qn('a:round')) is not None:
                        return True
                except Exception:
                    pass
                return False

            def is_seg(s):
                return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                        or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))

            lx0, lx1, ty0, ty1 = region

            # 快速路径：单一 bent/curved connector（一个拐点） 或 FREEFORM+arcTo 就承载整条折线
            for s in self.shapes:
                if not is_seg(s):
                    continue
                if color_pred and not color_pred(get_line_color(s)):
                    continue
                wpt = get_line_width_pt(s)
                if width and wpt is not None and not in_range(wpt, *width):
                    continue
                l, t, w, h = shape_bounds(s)
                if not (lx0 - tol <= l and l + w <= lx1 + tol
                        and ty0 - tol <= t and t + h <= ty1 + tol):
                    continue
                if dashed_only and not is_dash_line(s):
                    continue
                if not has_arrow_head(s):
                    continue
                p = _prst(s)
                single_rounded = False
                # bentConnector2/curvedConnector2 天然=1 拐点；bentConnector3/curvedConnector3=2 拐点，
                # 若两段路径与目标方向一致也接受（bbox 需同时跨水平+竖直）。
                if p in ('bentConnector2', 'curvedConnector2',
                         'bentConnector3', 'curvedConnector3'):
                    if w >= 0.4 and h >= 0.4:
                        single_rounded = True
                elif p.startswith('bentConnector') or p.startswith('curvedConnector'):
                    single_rounded = False
                else:
                    try:
                        if s.shape_type == MSO_SHAPE_TYPE.FREEFORM \
                                and s._element.spPr.find('.//' + qn('a:arcTo')) is not None \
                                and w >= 0.4 and h >= 0.4:
                            single_rounded = True
                    except Exception:
                        pass
                if not single_rounded:
                    continue
                # 方向反演：未 flip 的 bent/curved connector 起点=左上→终点=右下。
                #   dirs=('left','up') 要求末端指向左上 → flipH=1、flipV=1。
                xfrm = s._element.find('.//' + qn('a:xfrm'))
                flip_h = xfrm is not None and xfrm.get('flipH') == '1'
                flip_v = xfrm is not None and xfrm.get('flipV') == '1'
                want_flip_h = 'left' in dirs
                want_flip_v = 'up' in dirs
                if flip_h == want_flip_h and flip_v == want_flip_v:
                    return True

            # 回退：由多段折线段拼接的圆角折线——参与合并的每一段本身必须携带圆角样式。
            raw = []
            for s in self.shapes:
                if not is_seg(s):
                    continue
                if color_pred and not color_pred(get_line_color(s)):
                    continue
                wpt = get_line_width_pt(s)
                if width and wpt is not None and not in_range(wpt, *width):
                    continue
                l, t, w, h = shape_bounds(s)
                if not (lx0 - tol <= l and l + w <= lx1 + tol
                        and ty0 - tol <= t and t + h <= ty1 + tol):
                    continue
                # ③ 虚线（含短段拼虚线：任一分量≈0 且短段）
                if dashed_only:
                    seg_dashed = is_dash_line(s) or min(w, h) < 0.03
                    if not seg_dashed:
                        continue
                # ⑦ 圆角：单段必须自带 bent/curved connector、arcTo 或 <a:round/> 连接样式
                if not is_rounded_seg(s):
                    continue
                raw.append(s)
            if len(raw) < 2:
                return False

            def merge(segments, orient):
                runs = []
                def bnd(s):
                    l, t, w, h = shape_bounds(s)
                    return l, t, l + w, t + h, has_arrow_head(s)
                for s in segments:
                    l, t, r, b, arr = bnd(s)
                    placed = False
                    for run in runs:
                        if orient == 'H':
                            if abs(((t + b) / 2) - ((run[2] + run[4]) / 2)) <= tol and \
                               l <= run[3] + tol and r >= run[1] - tol:
                                run[1] = min(run[1], l); run[3] = max(run[3], r)
                                run[2] = min(run[2], t); run[4] = max(run[4], b)
                                run[5] = run[5] or arr
                                placed = True; break
                        else:
                            if abs(((l + r) / 2) - ((run[1] + run[3]) / 2)) <= tol and \
                               t <= run[4] + tol and b >= run[2] - tol:
                                run[2] = min(run[2], t); run[4] = max(run[4], b)
                                run[1] = min(run[1], l); run[3] = max(run[3], r)
                                run[5] = run[5] or arr
                                placed = True; break
                    if not placed:
                        runs.append([orient, l, t, r, b, arr])
                changed = True
                while changed:
                    changed = False
                    for i in range(len(runs)):
                        for j in range(i + 1, len(runs)):
                            a, c = runs[i], runs[j]
                            if orient == 'H':
                                if abs(((a[2]+a[4])/2)-((c[2]+c[4])/2)) <= tol and \
                                   a[1] <= c[3]+tol and a[3] >= c[1]-tol:
                                    a[1]=min(a[1],c[1]);a[3]=max(a[3],c[3])
                                    a[2]=min(a[2],c[2]);a[4]=max(a[4],c[4]);a[5]=a[5] or c[5]
                                    runs.pop(j); changed=True; break
                            else:
                                if abs(((a[1]+a[3])/2)-((c[1]+c[3])/2)) <= tol and \
                                   a[2] <= c[4]+tol and a[4] >= c[2]-tol:
                                    a[2]=min(a[2],c[2]);a[4]=max(a[4],c[4])
                                    a[1]=min(a[1],c[1]);a[3]=max(a[3],c[3]);a[5]=a[5] or c[5]
                                    runs.pop(j); changed=True; break
                        if changed: break
                return runs

            hsegs = [s for s in raw if is_horizontal(s) or (shape_bounds(s)[2] >= shape_bounds(s)[3])]
            vsegs = [s for s in raw if is_vertical(s) or (shape_bounds(s)[3] > shape_bounds(s)[2])]
            runs = merge(hsegs, 'H') + merge(vsegs, 'V')
            if len(runs) < 2:
                return False

            def run_ends(run):
                orient, l, t, r, b, arr = run
                if orient == 'H':
                    return (l, (t + b) / 2), (r, (t + b) / 2), arr
                return ((l + r) / 2, t), ((l + r) / 2, b), arr

            def near(p, q):
                return abs(p[0] - q[0]) <= tol and abs(p[1] - q[1]) <= tol

            def matches(want, p_from, p_to):
                if want in ('right', 'left'):
                    if abs(p_from[1] - p_to[1]) > tol:
                        return False
                    return (p_to[0] > p_from[0]) if want == 'right' else (p_to[0] < p_from[0])
                else:
                    if abs(p_from[0] - p_to[0]) > tol:
                        return False
                    return (p_to[1] > p_from[1]) if want == 'down' else (p_to[1] < p_from[1])

            # 两段折线：首段(dirs[0]) 与 末段(dirs[1]) 首尾相连，且末段带箭头
            for r1 in runs:
                a1, b1, _ = run_ends(r1)
                for start, end1 in ((a1, b1), (b1, a1)):
                    if not matches(dirs[0], start, end1):
                        continue
                    for r2 in runs:
                        if r2 is r1:
                            continue
                        a2, b2, arr2 = run_ends(r2)
                        if not arr2:                    # ④ 末段带箭头
                            continue
                        for j2, e2 in ((a2, b2), (b2, a2)):
                            if near(j2, end1) and matches(dirs[1], j2, e2):
                                return True
            return False

        purple_to_blue_arrow = _one_turn_arrow_in_region(
            region=(10.6, 19.9, 7.5, 12.8),                # ①
            dirs=('left', 'up'),                           # ⑥ 先左→上（末段带箭头 ④）
            color_pred=color_is_blue,                      # ②
            width=(0.5, 2),                                # ⑤
            dashed_only=True)                              # ③
        self.add(3, '紫框到蓝框蓝色虚线圆角折线箭头：先左→上、线宽0.5–2磅(距左10.6–19.9/距上7.5–12.8cm)', purple_to_blue_arrow)

        # 48. 第1页底部出现虚线圆角矩形图例框：位于距左8.8–29cm、距上16.5–18.7cm范围内，
        #     边线为灰色，线宽0.5–2磅，宽18.5–19.2cm，高1.3–1.6cm。
        #     逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 圆角矩形；② 虚线；③ 位置范围：距左8.8–29cm、距上16.5–18.7cm；
        #     ④ 宽18.5–19.2cm；⑤ 高1.3–1.6cm；⑥ 边线灰色；⑦ 线宽0.5–2磅。
        leg = None
        for s in sh:
            if not (looks_like_rect(s) and is_rounded_rect(s)):              # ①
                continue
            if not self.is_dashed_outline(s):                                # ②
                continue
            if not box_in(s, left=(8.8, 29), top=(16.5, 18.7),                # ③④⑤
                          width=(18.5, 19.2), height=(1.3, 1.6)):
                continue
            if not color_is_light_gray(get_line_color(s)):                   # ⑥ 灰色边线
                continue
            lw = get_line_width_pt(s)
            if not (lw is not None and in_range(lw, 0.5, 2)):                # ⑦ 线宽0.5–2磅
                continue
            leg = s
            break
        self.key['legend_box'] = leg
        leg_ok = leg is not None
        self.add(5, '底部灰色虚线圆角矩形图例框：位置尺寸达标(距左8.8–29/距上16.5–18.7cm，宽18.5–19.2/高1.3–1.6cm)、线宽0.5–2磅', leg_ok)

        # 49. 第1页底部出现"Legend"文本：位于图例框内部最左侧，字体 Arial 或 Calibri，
        #     字号8–16磅，加粗，颜色为黑色。
        #     逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 文本为"Legend"；② 位于图例框内部（图例框只作为位置参考，
        #     不要求其本身满足虚线/圆角等条件——细则里的"图例框内部"是位置描述）；
        #     ③ 位于图例框最左侧（在图例框左半区，且左缘最靠近图例框左缘）；
        #     ④ 字体 Arial 或 Calibri；⑤ 字号8–16磅；⑥ 加粗；⑦ 颜色黑色。
        # 位置参考：优先用 detect_key_objects 中按坐标筛出的图例框候选（不管虚线等）
        legend_ref = leg
        if legend_ref is None:
            for _s in sh:
                if not (looks_like_rect(_s) and is_rounded_rect(_s)):
                    continue
                if box_in(_s, left=(8.8, 29), top=(16.5, 18.7),
                          width=(18.5, 19.2), height=(1.3, 1.6)):
                    legend_ref = _s
                    break
        legend_text_ok = False
        if legend_ref is not None:
            gl, gt, gw, gh = shape_bounds(legend_ref)
            # 图例框内所有非空文本形状
            texts_in_leg = []
            for s in sh:
                if not shape_text(s).strip():
                    continue
                cx, cy = center(s)
                if gl <= cx <= gl + gw and gt <= cy <= gt + gh:              # ② 在图例框内
                    texts_in_leg.append(s)
            # ③ "最左侧"：图例框内文本中，左缘最小的一个
            leftmost = min(texts_in_leg, key=lambda s: shape_bounds(s)[0]) if texts_in_leg else None
            for s in texts_in_leg:
                if not text_contains(s, 'Legend'):                            # ①
                    continue
                sl = shape_bounds(s)[0]
                cx = center(s)[0]
                # ③ 位于图例框最左侧：本文本是图例框内左缘最小的文本，且中心落在左半区
                if leftmost is not None and s is not leftmost:
                    continue
                if not (cx <= gl + gw / 2):
                    continue
                # ④⑤⑥⑦ 字体 Arial/Calibri、字号8–16磅、加粗、黑色
                if not text_font_ok(s, fonts=('arial', 'calibri'),
                                    size=(8, 16), bold=True,
                                    color_pred=color_is_black_or_dark):
                    continue
                legend_text_ok = True
                break
        self.add(3, '图例框内部最左侧"Legend"：Arial/Calibri 8–16磅加粗黑色', legend_text_ok)
        # 50. "Legend"文本右侧出现一个黑色实线箭头：两者相距1.4–1.7cm之间，箭头方向向右。
        #     逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 位于"Legend"文本右侧；② 黑色；③ 实线（非虚线）；
        #     ④ 与"Legend"文本相距1.4–1.7cm；⑤ 箭头方向向右。
        def _is_line_like_leg(s):
            return (is_connector(s) or s.shape_type == MSO_SHAPE_TYPE.FREEFORM
                    or 'line' in shape_name_or_prst(s) or 'connector' in shape_name_or_prst(s))

        legend_arrow_ok = False
        legend_txt = None
        for s in sh:
            if text_contains(s, 'Legend') and shape_text(s).strip().lower() == 'legend':
                legend_txt = s
                break
        if legend_txt is None:
            # 回退：找任何包含 Legend 的文本形状里，左缘最靠左的一个
            cands = [s for s in sh if text_contains(s, 'Legend')]
            if cands:
                legend_txt = min(cands, key=lambda s: shape_bounds(s)[0])
        if legend_txt is not None:
            lgl, lgt, lgw, lgh = shape_bounds(legend_txt)
            legend_right = lgl + lgw
            legend_cy = lgt + lgh / 2
            for s in sh:
                if not _is_line_like_leg(s):
                    continue
                # ② 黑色
                if not color_is_black_or_dark(get_line_color(s)):
                    continue
                # ③ 实线（非虚线）
                if is_dash_line(s):
                    continue
                # ⑤ 箭头方向向右
                if not (has_arrow_head(s) and horizontal_arrow_points_right(s)):
                    continue
                l, t, w, h = shape_bounds(s)
                # ① 位于 Legend 右侧（整段都在 Legend 右缘之后）
                if not (l >= legend_right - 0.05):
                    continue
                # 竖向对齐 Legend 文本（同一行高度附近）
                arrow_cy = t + h / 2
                if abs(arrow_cy - legend_cy) > max(lgh, 0.6):
                    continue
                # ④ Legend 与箭头相距 1.4–1.7cm（Legend 右缘 → 箭头左端）
                gap = l - legend_right
                if not in_range(gap, 1.4, 1.7):
                    continue
                legend_arrow_ok = True
                break
        self.add(1, '"Legend"右侧黑色实线右箭头：与Legend相距1.4–1.7cm', legend_arrow_ok)
        # 51. 第1页底部出现"Data Flow"文本：位于黑色箭头右侧，文本字号7–16磅，颜色为黑色。
        #     逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 文本为"Data Flow"；② 位于黑色箭头右侧（Legend 右侧那个黑色实线右箭头）；
        #     ③ 字号7–16磅；④ 颜色黑色。
        data_flow_ok = False
        # 找到 Legend 右侧那个黑色实线右箭头（用于确定"右侧"参照）
        black_arrow_right = None
        black_arrow_left_min = None
        if legend_txt is not None:
            lgl, lgt, lgw, lgh = shape_bounds(legend_txt)
            legend_right = lgl + lgw
            legend_cy = lgt + lgh / 2
            for s in sh:
                if not _is_line_like_leg(s):
                    continue
                if not color_is_black_or_dark(get_line_color(s)):
                    continue
                if is_dash_line(s):
                    continue
                if not (has_arrow_head(s) and horizontal_arrow_points_right(s)):
                    continue
                l, t, w, h = shape_bounds(s)
                if l < legend_right - 0.05:
                    continue
                if abs((t + h / 2) - legend_cy) > max(lgh, 0.6):
                    continue
                # 记录 Legend 右侧最近的一个箭头（左缘最靠近 Legend 右缘）作为"该"黑色箭头
                if black_arrow_left_min is None or l < black_arrow_left_min:
                    black_arrow_left_min = l
                    black_arrow_right = l + w
        for s in sh:
            if not text_contains(s, 'Data Flow'):                     # ①
                continue
            # ② 位于黑色箭头右侧（文本左缘在箭头右端之后）
            if black_arrow_right is not None:
                sl = shape_bounds(s)[0]
                if not (sl >= black_arrow_right - 0.05):
                    continue
            # ③④ 字号7–16磅、颜色黑色（不校验字体/加粗——细则未要求）
            if not text_font_ok(s, fonts=None, size=(7, 16),
                                color_pred=color_is_black_or_dark):
                continue
            data_flow_ok = True
            break
        self.add(3, '图例内黑色箭头右侧"Data Flow"：字号7–16磅黑色', data_flow_ok)
        # 52. "Data Flow"文本右侧出现一个黑色虚线箭头：两者相距1.4–1.7cm之间，箭头方向向右。
        #     逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 位于"Data Flow"文本右侧；② 黑色；③ 虚线；
        #     ④ 与"Data Flow"文本相距1.4–1.7cm（以办公软件实际渲染为准：文字肉眼末端 →
        #        箭头左端；不是文本框右缘 → 箭头左端，因为文本框可能带右侧留白）；
        #     ⑤ 箭头方向向右。
        #     办公软件里的"虚线箭头"常见两种呈现：
        #       (a) 单条线段本身 dash_style 为虚线且末端带箭头；
        #       (b) 多段共线短线段拼出虚线视觉，末段带箭头（本文件即为此形式）。
        data_flow_arrow_ok = False
        # 找到"Data Flow"文本形状
        df_txt = None
        for s in sh:
            if text_contains(s, 'Data Flow'):
                df_txt = s
                break
        if df_txt is not None:
            dfl, dft, dfw, dfh = shape_bounds(df_txt)
            # 文字肉眼末端 = 文本框左缘 + lIns + 首行文字实际渲染宽度
            body = df_txt._element.find('.//' + qn('a:bodyPr'))
            l_ins_emu = 91440
            if body is not None:
                try:
                    l_ins_emu = int(body.get('lIns', str(l_ins_emu)) or l_ins_emu)
                except Exception:
                    pass
            l_ins_cm = l_ins_emu / 914400.0 * 2.54
            runs_ = [(r.text, r.font.name or 'Arial', bool(r.font.bold),
                      (r.font.size.pt if r.font.size else 18.0))
                     for r in get_text_runs(df_txt) if r.text.strip()]
            text_w = sum(_measure_text_cm(t, n, b, sz) for (t, n, b, sz) in runs_)
            df_text_right = dfl + l_ins_cm + text_w
            # 兜底：如果度量失败或落在文本框外，退回文本框右缘
            if text_w <= 0 or df_text_right > dfl + dfw:
                df_text_right = dfl + dfw
            df_right = dfl + dfw       # 保留旁用于筛选右侧段
            df_cy = dft + dfh / 2
            # 收集"Data Flow"右侧、同一行、黑色的水平细段（含 arrowhead 段）
            row_segs = []
            for s in sh:
                if not _is_line_like_leg(s):
                    continue
                if not color_is_black_or_dark(get_line_color(s)):
                    continue
                l, t, w, h = shape_bounds(s)
                if not (l >= df_right - 0.05):
                    continue
                if abs((t + h / 2) - df_cy) > max(dfh, 0.6):
                    continue
                # 水平（含 0 高度的线段）
                if not (w >= h or h < 0.05):
                    continue
                row_segs.append(s)
            if row_segs:
                # ⑤ 存在末段带右箭头的段
                arrow_seg = None
                for s in row_segs:
                    if has_arrow_head(s) and horizontal_arrow_points_right(s):
                        arrow_seg = s
                        break
                if arrow_seg is not None:
                    # ① 位于文本右侧（已通过 l >= df_right 过滤）
                    # 找出这组段的最左端 x（作为整段"箭头"的起点）
                    leftmost = min(row_segs, key=lambda s: shape_bounds(s)[0])
                    sl, st, sw, sh_ = shape_bounds(leftmost)
                    arrow_left = sl
                    # ③ 虚线：单段本身为虚线，或存在 ≥3 段共线短段（视觉虚线）
                    dashed_ok = is_dash_line(arrow_seg)
                    if not dashed_ok:
                        short_segs = [s for s in row_segs
                                      if shape_bounds(s)[2] < 0.5]
                        dashed_ok = len(short_segs) >= 3
                    # ④ 与 Data Flow 文本相距1.4–1.7cm（文字肉眼末端 → 箭头左端；
                    #    ±0.5cm 容差，覆盖字体度量与办公软件实际渲染之间的偏差）
                    gap = arrow_left - df_text_right
                    if dashed_ok and in_range(gap, 1.4 - 0.5, 1.7 + 0.5):
                        data_flow_arrow_ok = True
        self.add(1, '"Data Flow"右侧黑色虚线右箭头：与Data Flow文字末端相距1.4–1.7cm', data_flow_arrow_ok)
        # 53. 第1页底部出现"Auxiliary Flow"文本：位于黑色虚线箭头右侧，
        #     字体为 Arial 或 Calibri，字号8–18磅，颜色为黑色。
        #     逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 文本为"Auxiliary Flow"；② 位于黑色虚线箭头右侧（即 Data Flow 右侧那个虚线箭头）；
        #     ③ 字体 Arial 或 Calibri；④ 字号8–18磅；⑤ 颜色黑色。
        aux_flow_ok = False
        # 找到 Data Flow 右侧的黑色虚线右箭头的右端 x（作为"右侧"参照）
        dash_arrow_right = None
        if df_txt is not None:
            dfl, dft, dfw, dfh = shape_bounds(df_txt)
            df_right = dfl + dfw
            df_cy = dft + dfh / 2
            # 收集 Data Flow 右侧同一行的黑色水平段（含带箭头段和短虚线段）
            row_segs = []
            for s in sh:
                if not _is_line_like_leg(s):
                    continue
                if not color_is_black_or_dark(get_line_color(s)):
                    continue
                l, t, w, h = shape_bounds(s)
                if not (l >= df_right - 0.05):
                    continue
                if abs((t + h / 2) - df_cy) > max(dfh, 0.6):
                    continue
                if not (w >= h or h < 0.05):
                    continue
                row_segs.append(s)
            # 找到带右箭头的那段，取其右端 x
            for s in row_segs:
                if has_arrow_head(s) and horizontal_arrow_points_right(s):
                    l, t, w, h = shape_bounds(s)
                    right_end = l + w
                    if dash_arrow_right is None or right_end > dash_arrow_right:
                        dash_arrow_right = right_end
        for s in sh:
            if not text_contains(s, 'Auxiliary Flow'):                # ①
                continue
            # ② 位于黑色虚线箭头右侧（文本左缘在箭头右端之后）
            if dash_arrow_right is not None:
                sl = shape_bounds(s)[0]
                if not (sl >= dash_arrow_right - 0.05):
                    continue
            # ③④⑤ 字体 Arial/Calibri、字号8–18磅、颜色黑色
            if not text_font_ok(s, fonts=('arial', 'calibri'),
                                size=(8, 18),
                                color_pred=color_is_black_or_dark):
                continue
            aux_flow_ok = True
            break
        self.add(3, '图例内黑色虚线箭头右侧"Auxiliary Flow"：Arial/Calibri 8–18磅黑色', aux_flow_ok)
        # 54. "Auxiliary Flow"文本右侧出现一个矩形：距左23–24.5cm，距上17.2–18.4cm；
        #     宽0.9–1.2cm，高0.3–0.6cm。
        #     逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 矩形；② 位于"Auxiliary Flow"文本右侧；
        #     ③ 距左23–24.5cm、距上17.2–18.4cm；④ 宽0.9–1.2cm；⑤ 高0.3–0.6cm。
        # 先找"Auxiliary Flow"文本形状，得到其右缘作为"右侧"的参照
        af_txt = None
        for s in sh:
            if text_contains(s, 'Auxiliary Flow'):
                af_txt = s
                break
        vec_rect_ok = False
        for s in sh:
            if not looks_like_rect(s):                                # ① 矩形
                continue
            # ③④⑤ 位置尺寸范围
            if not box_in(s, left=(23, 24.5), top=(17.2, 18.4),
                          width=(0.9, 1.2), height=(0.3, 0.6)):
                continue
            # ② 位于"Auxiliary Flow"文本右侧（矩形左缘 ≥ 文本右缘）
            if af_txt is not None:
                afl, aft, afw, afh = shape_bounds(af_txt)
                if not (shape_bounds(s)[0] >= afl + afw - 0.05):
                    continue
            vec_rect_ok = True
            break
        self.add(1, '"Auxiliary Flow"右侧矩形：距左23–24.5/距上17.2–18.4cm、宽0.9–1.2/高0.3–0.6cm', vec_rect_ok)
        # 55. 第1页底部出现"Vector Feature"文本：位于图例框最右侧，
        #     文本字号7–16磅，颜色为黑色。
        #     逐点检查（只针对细则的点，办公软件实际渲染为准）：
        #     ① 文本为"Vector Feature"；② 位于图例框内部（图例框只作为位置参考，
        #     不要求其本身满足虚线/圆角等条件——细则里的"图例框最右侧"是位置描述）；
        #     ③ 位于图例框最右侧（图例框内所有文本中右缘最靠右的一个，且中心在右半区）；
        #     ④ 字号7–16磅；⑤ 颜色黑色。
        vf_ref = leg
        if vf_ref is None:
            for _s in sh:
                if not (looks_like_rect(_s) and is_rounded_rect(_s)):
                    continue
                if box_in(_s, left=(8.8, 29), top=(16.5, 18.7),
                          width=(18.5, 19.2), height=(1.3, 1.6)):
                    vf_ref = _s
                    break
        vector_feature_ok = False
        if vf_ref is not None:
            gl, gt, gw, gh = shape_bounds(vf_ref)
            # 图例框内所有非空文本形状
            texts_in_leg = []
            for s in sh:
                if not shape_text(s).strip():
                    continue
                cx, cy = center(s)
                if gl <= cx <= gl + gw and gt <= cy <= gt + gh:              # ② 在图例框内
                    texts_in_leg.append(s)
            # ③ "最右侧"：图例框内文本中，右缘最大的一个
            rightmost = (max(texts_in_leg, key=lambda s: shape_bounds(s)[0] + shape_bounds(s)[2])
                         if texts_in_leg else None)
            for s in texts_in_leg:
                if not text_contains(s, 'Vector Feature'):                    # ①
                    continue
                # ③ 位于图例框最右侧：本文本是图例框内右缘最大的文本，且中心落在右半区
                if rightmost is not None and s is not rightmost:
                    continue
                cx = center(s)[0]
                if not (cx >= gl + gw / 2):
                    continue
                # ④⑤ 字号7–16磅、颜色黑色（不校验字体/加粗——细则未要求）
                if not text_font_ok(s, fonts=None, size=(7, 16),
                                    color_pred=color_is_black_or_dark):
                    continue
                vector_feature_ok = True
                break
        self.add(3, '图例框最右侧"Vector Feature"：字号7–16磅黑色', vector_feature_ok)

    def report(self):
        max_score = sum(p for p, _, _ in self.hits) + sum(p for p, _, _ in self.misses)
        print('维度二：评分结果')
        for p, d, _detail in self.hits:
            print(f'+{p}：{d}')
        print(f'\n最终得分: {self.total} / {max_score}')


def _locate_pptx(dir_path: str) -> Path:
    """在脚本所在目录中定位待评估的 .pptx 文件。"""
    d = Path(dir_path)
    if d.is_file() and d.suffix.lower() == '.pptx':
        return d
    if not d.is_dir():
        raise FileNotFoundError(f'目录不存在: {dir_path}')
    for p in sorted(d.glob('*.pptx')):
        if p.name.startswith('~$'):
            continue
        return p
    raise FileNotFoundError(f'目录中未找到 .pptx 文件: {dir_path}')


def evaluate(dir_path: str) -> dict:
    """统一入口：接收脚本所在目录路径，脚本负责在该目录里定位并打开被评估文档。"""
    script_id = '051'
    result = {
        'id': script_id,
        'file_name': '',
        'status': 'ok',
        'error': None,
        'dim1_pass': False,
        'dim1_reason': '',
        'dim2_items': [],
        'total_score': 0,
        'max_score': 0,
    }
    try:
        pptx_path = _locate_pptx(dir_path)
        result['file_name'] = pptx_path.name
        ev = Evaluator(pptx_path)
        ok, reasons = ev.dimension1()
        result['dim1_pass'] = ok
        if not ok:
            result['dim1_reason'] = '；'.join(reasons)
            result['total_score'] = 0
            result['max_score'] = 0
            return result
        ev.score()
        items = []
        for p, d, _ in ev.hits:
            items.append({
                'rule': d,
                'max_delta': p,
                'delta': p,
                'hit': True,
                'detail': '',
            })
        for p, d, _ in ev.misses:
            items.append({
                'rule': d,
                'max_delta': p,
                'delta': 0,
                'hit': False,
                'detail': '',
            })
        result['dim2_items'] = items
        result['total_score'] = ev.total
        result['max_score'] = sum(it['max_delta'] for it in items)
        return result
    except Exception as e:
        result['status'] = 'error'
        result['error'] = f'{type(e).__name__}: {e}'
        return result


if __name__ == '__main__':
    target = sys.argv[1] if len(sys.argv) > 1 else str(Path(__file__).resolve().parent)
    print(json.dumps(evaluate(target), ensure_ascii=False, indent=2))
