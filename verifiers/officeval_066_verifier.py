# -*- coding: utf-8 -*-
"""
自动评估脚本：终稿_表格图表关联数据更新版_新增表格折线图.pptx
按照评分细则执行维度2各得分点/扣分点独立判断，命中则累加。

对外仅暴露 evaluate(dir_path: str) -> dict：
  - 入参为脚本所在目录路径；脚本自行在该目录内定位 .pptx 被评估文档
  - 返回结构化评分字典（见 §2.2）
"""

import os
import sys
import re
import json
import zipfile

from pptx import Presentation
from pptx.util import Emu

# 脚本编号（与文件名 officeval_066_verifier.py 中的编号一致）
SCRIPT_ID = "066"


def _locate_pptx(dir_path):
    """在给定目录内定位被评估的 .pptx 文件；返回绝对路径或 None。"""
    if not dir_path or not os.path.isdir(dir_path):
        return None
    candidates = []
    for name in os.listdir(dir_path):
        low = name.lower()
        if low.endswith(".pptx"):
            # 跳过 Office 临时文件（以 ~$ 开头）
            if name.startswith("~$"):
                continue
            candidates.append(os.path.join(dir_path, name))
    if not candidates:
        return None
    # 若存在多个，按文件名排序取第一个，保证可复现
    candidates.sort()
    return candidates[0]

# ====== 标准答案 ======
EXPECTED_HEADERS = [
    "年份",
    "智慧门店响应率（%）",
    "线上订单占比（%）",
    "自助结算覆盖率（%）",
    "绿色到店客流占比（%）",
    "低碳配送完成率（%）",
    "末端接驳便利度（%）",
    "运营协同指数（0–100）",
]
# 7 个指标系列名（不含"年份"列）
EXPECTED_SERIES = [h for h in EXPECTED_HEADERS[1:]]
# 不带单位的"短名"，用于宽松匹配
EXPECTED_SERIES_SHORT = [
    "智慧门店响应率",
    "线上订单占比",
    "自助结算覆盖率",
    "绿色到店客流占比",
    "低碳配送完成率",
    "末端接驳便利度",
    "运营协同指数",
]

EXPECTED_DATA = {
    "2016": [31.8, 18.9, 24.6, 27.3, 22.6, 33.4, 29.7],
    "2017": [34.9, 23.7, 29.8, 31.5, 26.1, 36.8, 34.1],
    "2018": [42.6, 27.9, 35.2, 34.4, 30.6, 39.6, 39.8],
    "2019": [45.7, 33.8, 40.9, 37.2, 35.9, 43.8, 44.7],
    "2020": [53.4, 41.7, 49.6, 42.8, 44.3, 50.6, 52.9],
    "2021": [58.6, 48.8, 55.7, 47.6, 49.7, 55.1, 58.4],
    "2022": [66.7, 56.9, 62.6, 55.8, 57.4, 63.7, 66.2],
    "2023": [71.6, 62.4, 69.8, 61.6, 64.9, 70.2, 72.1],
    "2024": [78.9, 68.1, 76.7, 67.3, 72.6, 74.9, 78.2],
}
EXPECTED_YEARS = list(EXPECTED_DATA.keys())

EXPECTED_TITLE_S5 = "门店数字化与低碳运营指标"
EXPECTED_TITLE_S6 = "门店数字化与低碳运营指标趋势"
EXPECTED_BG_HEX = "F7FAFC"


# ---------- 工具函数 ----------
def norm(s):
    """归一化：去空白、统一全角括号 / 短横线，便于宽松比较"""
    if s is None:
        return ""
    s = str(s)
    s = s.replace("\n", "").replace("\r", "").replace(" ", "").replace("　", "")
    s = s.replace("（", "(").replace("）", ")")
    s = s.replace("–", "-").replace("—", "-").replace("−", "-")
    return s


def approx_equal(a, b, tol=1e-6):
    try:
        return abs(float(a) - float(b)) <= tol
    except Exception:
        return False


def emu_to_pt(emu):
    if emu is None:
        return None
    return float(emu) / 12700.0


def _find_bg_hex_in_element(element):
    """在某个 XML 元素（slide / layout / master）下查找 <p:bg> 中的 srgbClr"""
    if element is None:
        return None
    a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    p_ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    bg_el = element.find(p_ns + "cSld/" + p_ns + "bg")
    if bg_el is None:
        return None
    for srgb in bg_el.iter(a_ns + "srgbClr"):
        return srgb.get("val", "").upper()
    return None


def get_slide_bg_hex(slide):
    """按 slide → layout → master 的顺序回溯背景填充色（仅识别纯色 srgb）"""
    # slide 自身
    h = _find_bg_hex_in_element(slide._element)
    if h:
        return h
    # slideLayout
    try:
        layout = slide.slide_layout
        h = _find_bg_hex_in_element(layout._element)
        if h:
            return h
        # slideMaster
        master = layout.slide_master
        h = _find_bg_hex_in_element(master._element)
        if h:
            return h
    except Exception:
        pass
    return None


def find_table(slide):
    for sh in slide.shapes:
        if sh.has_table:
            return sh
    return None


def find_chart(slide):
    for sh in slide.shapes:
        if sh.has_chart:
            return sh
    return None


def find_title_textbox(slide, expected_keywords=None):
    """找到位于页面左上方、字号 24 加粗的标题文本框；返回 (shape, text, font_name, size_pt, bold)"""
    candidates = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if not txt:
            continue
        # 收集第一段第一个 run 的字体信息
        font_name, size_pt, bold = None, None, None
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                font_name = r.font.name
                size_pt = emu_to_pt(r.font.size) if r.font.size else None
                bold = r.font.bold
                break
            if font_name or size_pt or bold is not None:
                break
        candidates.append({
            "shape": sh,
            "text": txt,
            "font": font_name,
            "size": size_pt,
            "bold": bold,
            "left": sh.left or 0,
            "top": sh.top or 0,
        })

    # 启发式：优先取（字号≈24 且 加粗）的；若无则取最靠左上的非空文本
    big_bold = [c for c in candidates if c["size"] and abs(c["size"] - 24) < 0.5 and c["bold"]]
    if big_bold:
        # 选择最靠上的
        return sorted(big_bold, key=lambda c: (c["top"], c["left"]))[0]
    if candidates:
        return sorted(candidates, key=lambda c: (c["top"], c["left"]))[0]
    return None


# ---------- 维度 2 ----------
def check_pages_order(prs):
    """+5：最终共 6 页，第 5 页为数据表页，第 6 页为对应折线图页。"""
    reasons = []

    # 1) 共 6 页
    n = len(prs.slides)
    if n != 6:
        return False, f"页数为 {n}，不等于 6"
    reasons.append("共 6 页")

    # 2) 第 5 页为数据表页：该页存在表格对象
    s5_table_shape = find_table(prs.slides[4])
    if s5_table_shape is None:
        return False, "第 5 页不是数据表页（未发现表格对象）"
    reasons.append("第 5 页为数据表页")

    # 3) 第 6 页为"对应"折线图页：该页存在折线图，且其类目/数据与第 5 页表格对应
    s6_chart_shape = find_chart(prs.slides[5])
    if s6_chart_shape is None:
        return False, "第 6 页不是折线图页（未发现图表对象）"
    chart = s6_chart_shape.chart
    ct = str(chart.chart_type).upper()
    if "LINE" not in ct:
        return False, f"第 6 页图表不是折线图（类型={ct}）"

    # "对应"判定：图表横轴类目集合 == 第 5 页表格首列（年份）集合
    table = s5_table_shape.table
    table_years = [norm(table.cell(r, 0).text) for r in range(1, len(table.rows))]
    chart_cats = []
    for plot in chart.plots:
        chart_cats = [norm(str(x)) for x in plot.categories]
        break
    if set(chart_cats) != set(table_years) or not chart_cats:
        return False, f"第 6 页折线图与第 5 页表格不对应：图表类目={chart_cats}, 表格年份={table_years}"
    reasons.append("第 6 页为对应折线图页")

    return True, "；".join(reasons)


def check_s5_title(prs):
    """+1：第5页标题
        细则要求逐条核验：
          (a) 使用"门店数字化与低碳运营指标"标题
          (b) 位于页面左上方顶部
          (c) 字体为微软雅黑
          (d) 24 磅
          (e) 加粗
    """
    s5 = prs.slides[4]
    info = find_title_textbox(s5)
    if not info:
        return False, "未找到第 5 页标题文本"

    slide_w = prs.slide_width or 9144000
    slide_h = prs.slide_height or 6858000

    # (a) 标题文本
    text_ok = norm(info["text"]) == norm(EXPECTED_TITLE_S5)
    # (b) 位于"页面左上方顶部"：水平在左半部分，垂直在顶部区域
    pos_ok = (info["left"] is not None and info["top"] is not None
              and info["left"] <= slide_w / 2
              and info["top"] <= slide_h / 3)
    # (c) 字体：微软雅黑
    font_ok = (info["font"] or "").replace(" ", "") in (
        "微软雅黑", "MicrosoftYaHei", "MicrosoftYaHeiUI"
    )
    # (d) 字号 24 磅
    size_ok = info["size"] is not None and abs(info["size"] - 24) < 0.5
    # (e) 加粗
    bold_ok = bool(info["bold"])

    checks = [
        ('标题文本="门店数字化与低碳运营指标"', text_ok),
        ("位于页面左上方顶部", pos_ok),
        ("字体为微软雅黑", font_ok),
        ("24 磅", size_ok),
        ("加粗", bold_ok),
    ]
    miss = [name for name, ok in checks if not ok]
    detail = (f"文本={info['text']!r} 字体={info['font']} 大小={info['size']} "
              f"加粗={info['bold']} 位置=(L={info['left']},T={info['top']})")
    if miss:
        detail += "；不达标：" + "，".join(miss)
        return False, detail
    return True, detail


def check_s5_table_structure(prs):
    """+5：第5页表格结构
        细则要求逐条核验：
          (a) 包含 1 行表头和 9 行数据
          (b) 共 10 行 8 列
          (c) 所有单元格可编辑
          (d) 表头从左至右依次为：年份 / 智慧门店响应率（%）/ 线上订单占比（%）/
              自助结算覆盖率（%）/ 绿色到店客流占比（%）/ 低碳配送完成率（%）/
              末端接驳便利度（%）/ 运营协同指数（0–100）
          (e) 左侧从上到下年份依次为 2016—2024
          (f) 各年份对应 7 个数值与细则完全一致
    """
    s5 = prs.slides[4]
    tsh = find_table(s5)
    if tsh is None:
        return False, "第 5 页未发现表格"
    t = tsh.table
    rows, cols = len(t.rows), len(t.columns)

    # (b) 共 10 行 8 列
    shape_ok = (rows == 10 and cols == 8)
    # (a) 1 行表头 + 9 行数据（等价于 rows==10，但单独列出以踩点）
    header_count_ok = (rows >= 1)
    data_row_count_ok = (rows - 1 == 9)

    # (c) 所有单元格可编辑：能访问每个单元格的 text_frame 即视为原生可编辑表格
    editable_ok = True
    if shape_ok:
        try:
            for r in range(rows):
                for c in range(cols):
                    _ = t.cell(r, c).text_frame
        except Exception:
            editable_ok = False

    # (d) 表头逐列严格匹配（归一化后字符串相等，去换行/空格、全角括号、短横线）
    header_actual = [norm(t.cell(0, c).text) for c in range(cols)] if shape_ok else []
    header_expect = [norm(h) for h in EXPECTED_HEADERS]
    header_ok = (header_actual == header_expect)

    # (e) 左侧年份依次为 2016—2024
    years_actual = [norm(t.cell(r, 0).text) for r in range(1, rows)] if shape_ok else []
    years_expect = [norm(y) for y in EXPECTED_YEARS]
    years_ok = (years_actual == years_expect)

    # (f) 数据逐格匹配
    data_ok = True
    err_cells = []
    if shape_ok and years_ok:
        for ri, year in enumerate(EXPECTED_YEARS, start=1):
            for ci in range(1, 8):
                actual_raw = t.cell(ri, ci).text.strip()
                try:
                    actual = float(actual_raw)
                except Exception:
                    data_ok = False
                    err_cells.append(f"[{ri},{ci}] 非数字 {actual_raw!r}")
                    continue
                if not approx_equal(actual, EXPECTED_DATA[year][ci - 1], tol=0.01):
                    data_ok = False
                    err_cells.append(f"[{ri},{ci}] {actual}!={EXPECTED_DATA[year][ci-1]}")
    else:
        data_ok = False

    checks = [
        ("包含 1 行表头", header_count_ok),
        ("包含 9 行数据", data_row_count_ok),
        ("共 10 行 8 列", shape_ok),
        ("所有单元格可编辑", editable_ok),
        ("表头从左至右依次正确", header_ok),
        ("左侧年份从上到下为 2016—2024", years_ok),
        ("各年份 7 个数值与细则一致", data_ok),
    ]
    miss = [name for name, ok in checks if not ok]
    detail = (f"行列={rows}x{cols}, 表头一致={header_ok}, 年份一致={years_ok}, "
              f"数据一致={data_ok}, 可编辑={editable_ok}")
    if not header_ok:
        detail += f"; 表头实际={header_actual}"
    if not years_ok:
        detail += f"; 年份实际={years_actual}"
    if err_cells:
        detail += f"; 数据不一致样例={err_cells[:3]}"
    if miss:
        detail += "；不达标：" + "，".join(miss)
        return False, detail
    return True, detail


def _hex_close(hex_str, target_hex, max_dist=80):
    """两个颜色十六进制的曼哈顿距离判定"""
    if not hex_str or len(hex_str) < 6:
        return False
    hex_str = hex_str.upper()[-6:]
    target_hex = target_hex.upper()[-6:]
    try:
        r1, g1, b1 = int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16)
        r2, g2, b2 = int(target_hex[0:2], 16), int(target_hex[2:4], 16), int(target_hex[4:6], 16)
        return abs(r1 - r2) + abs(g1 - g2) + abs(b1 - b2) <= max_dist
    except Exception:
        return False


def _cell_fill_hex(cell):
    """取出单元格的填充色 hex；非纯色填充返回 None"""
    try:
        fill = cell.fill
        if fill.type == 1:  # SOLID
            return str(fill.fore_color.rgb).upper()
    except Exception:
        pass
    # 退而求其次：从 XML 解析
    tc = cell._tc
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for srgb in tc.iter(ns + "srgbClr"):
        return srgb.get("val", "").upper()
    return None


def _run_color_hex(cell):
    tc = cell._tc
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for rPr in tc.iter(ns + "rPr"):
        for srgb in rPr.iter(ns + "srgbClr"):
            return srgb.get("val", "").upper()
    return None


def check_s5_table_style(prs):
    """+1：第5页表格样式
        细则要求逐条核验：
          (a) 使用蓝色或青绿色表头（所有表头单元格底色均需符合）
          (b) 白色表头文字（所有表头单元格的文字颜色显式或解析后为白色）
          (c) 浅色交替数据行（每行所有数据单元格底色一致，且相邻行底色不同，
              两种底色均为浅色）
    """
    s5 = prs.slides[4]
    tsh = find_table(s5)
    if tsh is None:
        return False, "第 5 页未发现表格"
    t = tsh.table
    cols = len(t.columns)
    rows = len(t.rows)

    # (a) 表头底色：所有表头单元格均为蓝色或青绿色
    header_fills = [_cell_fill_hex(t.cell(0, c)) for c in range(cols)]

    def is_blue_or_teal(h):
        if not h or len(h) < 6:
            return False
        h = h.upper()[-6:]
        try:
            r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        except Exception:
            return False
        # 蓝色：B 显著大于 R 和 G；青绿色：G 与 B 都显著大于 R
        is_blue = (b > r + 20) and (b > g - 30) and (r < 200)
        is_teal = (g > r + 20) and (b > r + 20) and (r < 200)
        return is_blue or is_teal

    # 严格：所有表头单元格底色都必须显式命中蓝/青绿
    header_color_ok = (
        len(header_fills) == cols
        and cols > 0
        and all(is_blue_or_teal(h) for h in header_fills)
    )

    # (b) 白色表头文字：所有表头单元格的文字颜色必须显式或解析后为白色
    header_text_colors = []
    header_text_white = cols > 0
    for c in range(cols):
        rgb = _run_color_hex(t.cell(0, c))
        header_text_colors.append(rgb)
        # 未显式设置颜色 → 视为不合格（严格要求所有表头文字为白色）
        if rgb is None or not _hex_close(rgb, "FFFFFF", max_dist=60):
            header_text_white = False

    # (c) 浅色交替数据行：
    #     - 每行的所有数据单元格底色需一致（作为该行代表色）
    #     - 相邻数据行底色不同（形成"交替"）
    #     - 出现的两种底色均为浅色
    def is_light(h):
        if not h or len(h) < 6:
            return False
        h = h.upper()[-6:]
        try:
            r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        except Exception:
            return False
        # 浅色：RGB 三通道平均亮度 >= 200
        return (r + g + b) / 3 >= 200

    row_reps = []       # 每行的代表底色（None 表示该行内部不一致或缺失）
    row_uniform_ok = True
    for r in range(1, rows):
        row_fills = [_cell_fill_hex(t.cell(r, c)) for c in range(cols)]
        # 该行所有数据单元格必须都有底色且完全一致
        if any(x is None for x in row_fills) or len(set(row_fills)) != 1:
            row_uniform_ok = False
            row_reps.append(None)
        else:
            row_reps.append(row_fills[0])

    if row_uniform_ok and row_reps:
        distinct_set = set(row_reps)
        alt_two_colors_ok = len(distinct_set) >= 2
        alt_pattern_ok = True
        for i in range(len(row_reps) - 1):
            if row_reps[i] == row_reps[i + 1]:
                alt_pattern_ok = False
                break
        alt_light_ok = alt_two_colors_ok and all(is_light(x) for x in distinct_set)
        alt_ok = alt_two_colors_ok and alt_pattern_ok and alt_light_ok
    else:
        alt_ok = False

    checks = [
        ("蓝色或青绿色表头（所有表头单元格）", header_color_ok),
        ("白色表头文字（所有表头单元格）", header_text_white),
        ("浅色交替数据行（整行一致且相邻行不同）", alt_ok),
    ]
    miss = [name for name, ok in checks if not ok]
    detail = (f"表头底色={header_fills} 表头字色={header_text_colors} "
              f"数据行代表底色={row_reps} 行内一致={row_uniform_ok}")
    if miss:
        detail += "；不达标：" + "，".join(miss)
        return False, detail
    return True, detail


def _get_chart_embedded_workbook_data(prs, chart_shape):
    """读取折线图所嵌入工作簿（xlsx）中的实际数据。
    返回 dict:
       {
         'headers': [表头列表, 第一项通常是"年份"，后续是7个指标名],
         'data':    {年份(归一化): [7个数值]},
       }
    读取失败返回 None。
    """
    try:
        import zipfile, io
        chart_part = chart_shape.chart.part
        # 找到与 chart part 相关的 embedded xlsx
        wb_part = None
        for rel in chart_part.rels.values():
            if "spreadsheetml" in rel.reltype.lower() or rel.target_ref.lower().endswith(".xlsx"):
                wb_part = rel.target_part
                break
        if wb_part is None:
            return None
        xlsx_bytes = wb_part.blob
        try:
            from openpyxl import load_workbook
        except Exception:
            return None
        wb = load_workbook(io.BytesIO(xlsx_bytes), data_only=True)
        ws = wb.active
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            return None
        # 第一行是表头
        header_row = rows[0] if rows else ()
        headers = [("" if v is None else str(v)) for v in header_row]
        # 第一列是年份，后续构造 {年份: [7个数值]}
        data = {}
        for row in rows[1:]:
            if not row or row[0] is None:
                continue
            key = str(row[0]).strip()
            vals = [v for v in row[1:8]]
            if len(vals) < 7:
                continue
            data[norm(key)] = vals
        return {"headers": headers, "data": data}
    except Exception:
        return None


def check_s6_chart_data(prs):
    """+5：第6页折线图对象
        细则要求逐条核验：
          (a) 使用 PowerPoint 原生可编辑折线图
          (b) 图表数据工作簿中包含第 5 页全部年份和指标（比对工作簿表头 7 个指标名）
          (c) 横轴从左至右依次显示 2016、2017、…、2024
          (d) 数据系列完整包含 7 条折线（指定的 7 个名称）
          (e) 7 条折线的全部 63 个数据点与第 5 页对应单元格完全一致
          (f) 纵轴采用 0 至 100 的统一刻度（min/max 必须显式为 0 / 100）
          (g) 图例准确显示 7 个系列名称
          (h) 图例顺序与第 5 页表头顺序一致
          (i) 不使用"系列1"等默认名称
          (j) 图像数据来源于第五页表格（直接比对图表工作簿/系列数据与第 5 页实际单元格）
    """
    s6 = prs.slides[5]
    csh = find_chart(s6)
    if csh is None:
        return False, "第 6 页无图表"
    chart = csh.chart

    # (a) 原生可编辑折线图：chart_type 含 LINE 且 chart.part 能解析
    ct = str(chart.chart_type).upper()
    native_line_ok = ("LINE" in ct) and (chart.part is not None)

    # 解析图表中的系列与类目
    cats = []
    series_list = []
    for plot in chart.plots:
        cats = [str(x) for x in plot.categories]
        for s in plot.series:
            series_list.append((s.name, [float(v) if v is not None else None for v in s.values]))
        break

    # (c) 横轴从左至右依次 2016—2024
    cats_ok = cats == EXPECTED_YEARS

    # (d) 7 条折线 + 名称完整包含指定 7 个（去括号去单位后比较）
    def short(name):
        n = norm(name)
        n = re.sub(r"\(.*?\)", "", n)
        return n
    actual_short = [short(s[0]) for s in series_list]
    expected_short = [norm(x) for x in EXPECTED_SERIES_SHORT]
    n_series_ok = len(series_list) == 7
    series_names_ok = n_series_ok and set(actual_short) == set(expected_short)

    # ---- 直接从第 5 页表格读取实际单元格数据（作为"来源"真值）----
    s5 = prs.slides[4]
    s5_tsh = find_table(s5)
    s5_years = []            # 第 5 页表格中的年份（保持行顺序）
    s5_row_by_year = {}      # {年份(归一化): [7个数值 or None]}
    s5_metric_headers = []   # 第 5 页表头中除"年份"外的 7 个指标名
    if s5_tsh is not None:
        st = s5_tsh.table
        s5_rows_n = len(st.rows)
        s5_cols_n = len(st.columns)
        if s5_cols_n >= 8:
            s5_metric_headers = [st.cell(0, c).text.strip() for c in range(1, 8)]
        for r in range(1, s5_rows_n):
            y_raw = st.cell(r, 0).text.strip()
            if not y_raw:
                continue
            s5_years.append(y_raw)
            row_vals = []
            for c in range(1, min(8, s5_cols_n)):
                cell_txt = st.cell(r, c).text.strip()
                try:
                    row_vals.append(float(cell_txt))
                except Exception:
                    row_vals.append(None)
            # 补足 7 列
            while len(row_vals) < 7:
                row_vals.append(None)
            s5_row_by_year[norm(y_raw)] = row_vals

    # (e) 63 个数据点与第 5 页对应单元格完全一致
    #     以第 5 页表格实际数据为真值（而非固定的 EXPECTED_DATA），实现"来源"验证
    data_ok = True
    bad = []
    if n_series_ok and cats_ok and series_names_ok and s5_row_by_year:
        # 系列名 → 对应第 5 页表头中的列索引（0..6）
        s5_short_headers = [short(h) for h in s5_metric_headers]
        name_to_idx = {}
        for i, sh_name in enumerate(s5_short_headers):
            name_to_idx[sh_name] = i
        for ser_name, vals in series_list:
            col_idx = name_to_idx.get(short(ser_name))
            if col_idx is None:
                data_ok = False
                bad.append(f"系列名在第5页表头中未匹配: {ser_name}")
                continue
            for yi, year in enumerate(cats):
                row = s5_row_by_year.get(norm(year))
                if row is None:
                    data_ok = False
                    bad.append(f"第5页无年份 {year}")
                    break
                exp = row[col_idx]
                if exp is None or vals[yi] is None or not approx_equal(vals[yi], exp, tol=0.01):
                    data_ok = False
                    bad.append(f"[{ser_name} / {year}] {vals[yi]}!={exp}")
    else:
        data_ok = False

    # (f) 纵轴 0~100 的统一刻度：min 必须显式为 0，max 必须显式为 100（None 不再放行）
    va_min = None
    va_max = None
    try:
        va = chart.value_axis
        va_min = va.minimum_scale
        va_max = va.maximum_scale
        min_ok = va_min is not None and abs(va_min - 0) < 1e-6
        max_ok = va_max is not None and abs(va_max - 100) < 1e-6
        val_axis_ok = min_ok and max_ok
    except Exception:
        val_axis_ok = False

    # (g) 图例存在并显示 7 个系列名称
    legend_ok = chart.has_legend
    legend_count_ok = legend_ok and (len(series_list) == 7)

    # (h) 图例顺序与第 5 页表头顺序一致
    #     顺序真值取第 5 页实际表头（若不可用，退化为 EXPECTED_SERIES_SHORT）
    if s5_metric_headers and len(s5_metric_headers) == 7:
        header_order = [short(h) for h in s5_metric_headers]
    else:
        header_order = expected_short
    legend_order_ok = actual_short == header_order

    # (i) 不使用"系列N / SeriesN"等默认名
    no_default_name = all(
        not re.match(r"^系列\s*\d+$|^Series\s*\d+$", (s[0] or "").strip(), re.I)
        for s in series_list
    )

    # (b) 图表数据工作簿中包含第 5 页全部年份和指标 —— 比对工作簿表头 7 个指标名
    # (j) 图像数据来源于第五页表格 —— 直接比对工作簿数据与第 5 页实际单元格
    wb_info = _get_chart_embedded_workbook_data(prs, csh)
    wb_data = wb_info["data"] if wb_info else None
    wb_headers = wb_info["headers"] if wb_info else None
    wb_years_ok = False
    wb_metrics_ok = False
    wb_source_ok = False
    if wb_info is not None and wb_data is not None:
        # 年份覆盖：第 5 页全部年份都需在工作簿出现
        s5_year_keys = [norm(y) for y in s5_years] if s5_years else [norm(y) for y in EXPECTED_YEARS]
        wb_years_ok = bool(s5_year_keys) and all(y in wb_data for y in s5_year_keys)
        # 指标名比对：工作簿表头（去首列年份后）应与第 5 页表头 7 个指标名一致（去括号/空白后）
        wb_metric_headers = wb_headers[1:8] if wb_headers and len(wb_headers) >= 8 else []
        wb_metric_short = [short(h) for h in wb_metric_headers]
        s5_metric_short = [short(h) for h in s5_metric_headers] if s5_metric_headers else expected_short
        wb_metrics_ok = (
            len(wb_metric_short) == 7
            and len(s5_metric_short) == 7
            and wb_metric_short == s5_metric_short
        )
        # 数值来源：工作簿数据须与第 5 页实际单元格完全一致
        if wb_years_ok and wb_metrics_ok and s5_row_by_year:
            wb_source_ok = True
            for y in s5_year_keys:
                wb_row = wb_data.get(y)
                s5_row = s5_row_by_year.get(y)
                if wb_row is None or s5_row is None:
                    wb_source_ok = False
                    break
                for ci in range(7):
                    try:
                        if s5_row[ci] is None or not approx_equal(float(wb_row[ci]), s5_row[ci], tol=0.01):
                            wb_source_ok = False
                            break
                    except Exception:
                        wb_source_ok = False
                        break
                if not wb_source_ok:
                    break
    # 若无法读取嵌入工作簿：三项均判为不通过（不再退化，避免误放行）

    checks = [
        ("使用 PowerPoint 原生可编辑折线图", native_line_ok),
        ("图表数据工作簿包含第5页全部年份", wb_years_ok),
        ("图表数据工作簿包含第5页全部指标名", wb_metrics_ok),
        ("横轴依次为 2016—2024", cats_ok),
        ("数据系列完整包含 7 条折线（指定名称）", series_names_ok),
        ("63 个数据点与第 5 页完全一致", data_ok),
        ("纵轴显式 0—100 统一刻度", val_axis_ok),
        ("图例准确显示 7 个系列名称", legend_count_ok),
        ("图例顺序与第 5 页表头顺序一致", legend_order_ok),
        ("不使用'系列1'等默认名称", no_default_name),
        ("图像数据来源于第五页表格", wb_source_ok),
    ]
    miss = [name for name, ok in checks if not ok]
    detail = (f"类型={ct}, 横轴={cats}, 系列顺序={actual_short}, 数据一致={data_ok}, "
              f"纵轴min={va_min} 纵轴max={va_max}, "
              f"有图例={legend_ok}, 工作簿可读={wb_info is not None}, "
              f"工作簿表头={wb_headers if wb_info else None}, "
              f"第5页表头指标={s5_metric_headers}")
    if bad:
        detail += f"; 数据不符样例={bad[:3]}"
    if miss:
        detail += "；不达标：" + "，".join(miss)
        return False, detail
    return True, detail


def check_s6_line_style(prs):
    """+1：第6页折线样式
        细则要求逐条核验：
          (a) 各系列使用可区分的不同颜色线条：每个系列都必须能解析到"显式"颜色
              （srgbClr / schemeClr 中任一），且颜色两两不同
          (b) 设置清晰的数据点标记：每个系列都必须显式设置 marker，symbol 不为 "none"，
              size 若显式则 >= 5
    """
    s6 = prs.slides[5]
    csh = find_chart(s6)
    if csh is None:
        return False, "第 6 页无图表"
    chart = csh.chart
    ns_c = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    root = chart._chartSpace

    # 逐系列收集：线条颜色（srgb 或 scheme 标识）、marker 符号、marker 大小
    series_count = 0
    line_colors = []       # 每系列一个字符串标识（"srgb:XXXXXX" 或 "scheme:accent1" 等），无法解析为 None
    marker_symbols = []
    marker_sizes = []
    for ser in root.iter(ns_c + "ser"):
        series_count += 1
        # 线条颜色：ser/spPr/ln 下的 solidFill 中 srgbClr 或 schemeClr
        color_id = None
        spPr = ser.find(ns_c + "spPr")
        if spPr is not None:
            ln = spPr.find(ns_a + "ln")
            if ln is not None:
                # 优先取 solidFill 内的颜色，其次任意后代
                srgb_el = None
                for el in ln.iter(ns_a + "srgbClr"):
                    srgb_el = el
                    break
                if srgb_el is not None:
                    color_id = "srgb:" + (srgb_el.get("val") or "").upper()
                else:
                    scheme_el = None
                    for el in ln.iter(ns_a + "schemeClr"):
                        scheme_el = el
                        break
                    if scheme_el is not None:
                        color_id = "scheme:" + (scheme_el.get("val") or "").lower()
        line_colors.append(color_id)

        # 数据点标记
        mk = ser.find(ns_c + "marker")
        sym = None
        size = None
        if mk is not None:
            sym_el = mk.find(ns_c + "symbol")
            if sym_el is not None:
                sym = sym_el.get("val")
            size_el = mk.find(ns_c + "size")
            if size_el is not None:
                try:
                    size = int(size_el.get("val"))
                except Exception:
                    size = None
        marker_symbols.append(sym)
        marker_sizes.append(size)

    # (a) 各系列使用可区分的不同颜色线条
    #     严格：每个系列都必须能解析到显式颜色，且颜色两两不同
    if series_count == 0:
        color_distinct_ok = False
    else:
        all_explicit = all(c is not None for c in line_colors)
        color_distinct_ok = all_explicit and len(set(line_colors)) == series_count

    # (b) 设置清晰的数据点标记
    #     严格：每个系列都必须显式设置 marker，symbol 非空、非 "none"；
    #     size 若显式设置则要求 >= 5；不再依赖 chart_type 中的 "MARKER" 兜底
    if series_count == 0:
        marker_clear_ok = False
    else:
        marker_clear_ok = True
        for i in range(series_count):
            sym = marker_symbols[i]
            size = marker_sizes[i]
            if sym is None or sym == "none" or sym == "":
                marker_clear_ok = False
                break
            if size is not None and size < 5:
                marker_clear_ok = False
                break

    checks = [
        ("各系列使用可区分的不同颜色线条", color_distinct_ok),
        ("设置清晰的数据点标记", marker_clear_ok),
    ]
    miss = [name for name, ok in checks if not ok]
    detail = (f"系列数={series_count} 线条色={line_colors} "
              f"标记符号={marker_symbols} 标记大小={marker_sizes} "
              f"图类型={chart.chart_type}")
    if miss:
        detail += "；不达标：" + "，".join(miss)
        return False, detail
    return True, detail


def check_s6_title(prs):
    """+1：第6页图表标题与标注
        细则要求逐条核验：
          (a) 出现"标题"对象：chart 原生 title，或位于图表上方/页面左上方顶部、
              内容与预期标题相符（或至少含关键词）的独立文本框；
              普通正文（长段落、位于图表下方等）不视为标题
          (b) 位于页面左上方顶部
          (c) 字体为微软雅黑
          (d) 24 磅
          (e) 加粗
          (f) 图表下方有"最高点说明"：需位于图表下方，且描述图表中的最高数据点
              （出现"最高/最大/峰值/顶点"等语义关键词 + 全表最大值 78.9，
               并同时指向该点的系列名与年份）
    """
    s6 = prs.slides[5]
    slide_w = prs.slide_width or 9144000
    slide_h = prs.slide_height or 6858000

    csh = find_chart(s6)
    chart_top = chart_bottom = None
    if csh is not None:
        try:
            chart_top = csh.top or 0
            chart_bottom = chart_top + (csh.height or 0)
        except Exception:
            pass

    expect = norm(EXPECTED_TITLE_S6)

    # ---------- (a)~(e) 标题识别 ----------
    # 只把符合下述条件之一的对象视为"标题"：
    #   1) chart 原生 title；
    #   2) 独立文本框，且满足：文本较短（<= 40 字符）、单段落、位于图表上方或页面上 1/3 区域、
    #      文本与预期标题归一化相等，或至少包含"门店数字化"/"低碳运营"/"趋势"等关键词。
    chart_title_text = None
    if csh is not None and csh.chart.has_title:
        try:
            chart_title_text = csh.chart.chart_title.text_frame.text.strip()
        except Exception:
            chart_title_text = None
    chart_title_hit = bool(chart_title_text)

    # 严格挑选候选 TextBox
    title_keywords = ("门店数字化", "低碳运营", "运营指标", "趋势")

    def looks_like_title(shape):
        if not shape.has_text_frame:
            return False, None
        text = shape.text_frame.text.strip()
        if not text:
            return False, None
        # 排除长段落 / 多行段落（正文特征）
        paragraphs = [p for p in shape.text_frame.paragraphs if p.text.strip()]
        if len(paragraphs) > 1:
            return False, None
        if len(text) > 40:
            return False, None
        # 位置：必须在页面顶部 1/3；若图表已定位，则也应在图表上方
        top = shape.top or 0
        if top > slide_h / 3:
            return False, None
        if chart_top is not None and top >= chart_top:
            return False, None
        # 内容：与预期标题相等，或包含关键词
        n = norm(text)
        if n == expect or any(k in text for k in title_keywords):
            return True, text
        return False, None

    tb = None            # 命中的标题文本框描述
    tb_text = None
    for sh in s6.shapes:
        ok, text = looks_like_title(sh)
        if not ok:
            continue
        # 收集字体/字号/加粗（取第一段第一 run）
        font_name, size_pt, bold = None, None, None
        for para in sh.text_frame.paragraphs:
            for r in para.runs:
                font_name = r.font.name
                size_pt = emu_to_pt(r.font.size) if r.font.size else None
                bold = r.font.bold
                break
            if font_name or size_pt is not None or bold is not None:
                break
        tb = {
            "shape": sh,
            "text": text,
            "font": font_name,
            "size": size_pt,
            "bold": bold,
            "left": sh.left or 0,
            "top": sh.top or 0,
        }
        tb_text = text
        break

    tb_title_hit = tb is not None and bool(tb_text)

    # (a) 出现标题：chart 原生 title 或严格筛选的 TextBox
    text_ok = chart_title_hit or tb_title_hit

    # (b)~(e) 位置/字体/字号/加粗：以承载标题文本的对象为准；TextBox 优先
    pos_ok = font_ok = size_ok = bold_ok = False
    fmt_source = None
    if tb_title_hit and tb is not None:
        pos_ok = (tb["left"] <= slide_w / 2 and tb["top"] <= slide_h / 3)
        font_ok = (tb["font"] or "").replace(" ", "") in (
            "微软雅黑", "MicrosoftYaHei", "MicrosoftYaHeiUI"
        )
        size_ok = tb["size"] is not None and abs(tb["size"] - 24) < 0.5
        bold_ok = bool(tb["bold"])
        fmt_source = (f"TextBox 文本={tb['text']!r} 字体={tb['font']} 大小={tb['size']} "
                      f"加粗={tb['bold']} 位置=(L={tb['left']},T={tb['top']})")
    elif chart_title_hit and csh is not None:
        cleft, ctop = csh.left, csh.top
        pos_ok = (cleft is not None and ctop is not None
                  and cleft <= slide_w / 2
                  and ctop <= slide_h / 3)
        try:
            tf = csh.chart.chart_title.text_frame
            ft_name, ft_size, ft_bold = None, None, None
            for para in tf.paragraphs:
                for r in para.runs:
                    ft_name = r.font.name
                    ft_size = emu_to_pt(r.font.size) if r.font.size else None
                    ft_bold = r.font.bold
                    break
                if ft_name or ft_size is not None or ft_bold is not None:
                    break
            font_ok = (ft_name or "").replace(" ", "") in (
                "微软雅黑", "MicrosoftYaHei", "MicrosoftYaHeiUI"
            )
            size_ok = ft_size is not None and abs(ft_size - 24) < 0.5
            bold_ok = bool(ft_bold)
            fmt_source = (f"chart.title 文本={chart_title_text!r} 字体={ft_name} 大小={ft_size} "
                          f"加粗={ft_bold} 图表位置=(L={cleft},T={ctop})")
        except Exception:
            fmt_source = f"chart.title 文本={chart_title_text!r} 但无法读取格式"

    # ---------- (f) 最高点说明 ----------
    # 计算图表中的真实最高点（真值来自 EXPECTED_DATA），用于比对文本是否指向该点
    max_val = None
    max_year = None
    max_metric = None
    for year, vals in EXPECTED_DATA.items():
        for i, v in enumerate(vals):
            if max_val is None or v > max_val:
                max_val = v
                max_year = year
                max_metric = EXPECTED_SERIES_SHORT[i]
    max_val_str = f"{max_val:g}"  # e.g. "78.9"

    # 语义关键词：必须表达"最高/最大"这一概念，纯 78.9 不足以证明是"说明"
    semantic_keywords = ("最高", "最大", "峰值", "顶点", "最大值", "最高点")

    note_ok = False
    note_text_found = None
    for sh in s6.shapes:
        if not sh.has_text_frame:
            continue
        tx = sh.text_frame.text.strip()
        if not tx:
            continue
        # 排除标题本身
        if norm(tx) == expect:
            continue
        # 排除已识别的标题 TextBox
        if tb is not None and sh is tb["shape"]:
            continue
        # 必须位于图表下方
        sh_top = sh.top or 0
        if chart_bottom is None or sh_top < chart_bottom:
            continue
        # 必须同时满足：出现"最高/最大"语义关键词 + 指向真实最高点
        has_semantic = any(k in tx for k in semantic_keywords)
        # 指向真实最高点的最低条件：文本包含最大值（78.9），并至少同时提及该点的年份或指标名
        points_to_max = (
            (max_val_str in tx)
            and ((max_year and str(max_year) in tx) or (max_metric and max_metric in tx))
        )
        if has_semantic and points_to_max:
            note_ok = True
            note_text_found = tx
            break

    checks = [
        ("出现标题", text_ok),
        ("位于页面左上方顶部", pos_ok),
        ("字体为微软雅黑", font_ok),
        ("24 磅", size_ok),
        ("加粗", bold_ok),
        ("图表下方有最高点说明内容（指向真实最高点）", note_ok),
    ]
    miss = [name for name, ok in checks if not ok]

    detail = (f"{fmt_source}; 图表底={chart_bottom}; "
              f"最高点真值={max_metric}/{max_year}={max_val_str}; "
              f"最高点说明命中={note_ok}")
    if note_text_found:
        detail += f"(说明文本={note_text_found!r})"
    if miss:
        detail += "；不达标：" + "，".join(miss)
        return False, detail
    return True, detail


def _get_slide_bg_fill_info(slide):
    """读取一页幻灯片的"有效"背景填充信息，返回 dict:
       { 'type': 'solid'|'gradient'|'pattern'|'blip'|'ref'|None,
         'hex': 'FFFFFF' or None,
         'source': 'slide'|'layout'|'master'|None }
       按 slide → slideLayout → slideMaster 的顺序回溯：先看 slide 自身 <p:bg>，
       未显式设置则回溯到 layout / master，从而识别通过母版/版式设置的相同背景。
    """
    a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    p_ns = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

    def parse_bg_from_element(element):
        """从 slide/layout/master 元素中解析 <p:bg>；未显式设置返回 None。"""
        if element is None:
            return None
        bg_el = element.find(p_ns + "cSld/" + p_ns + "bg")
        if bg_el is None:
            return None
        bgPr = bg_el.find(p_ns + "bgPr")
        bgRef = bg_el.find(p_ns + "bgRef")
        if bgRef is not None:
            # <p:bgRef> 通过主题引用背景色；尝试取其内部的 srgbClr
            hex_val = None
            for srgb in bgRef.iter(a_ns + "srgbClr"):
                hex_val = (srgb.get("val") or "").upper()
                break
            return {"type": "ref", "hex": hex_val}
        if bgPr is None:
            return None
        if bgPr.find(a_ns + "solidFill") is not None:
            srgb = bgPr.find(a_ns + "solidFill/" + a_ns + "srgbClr")
            hex_val = srgb.get("val", "").upper() if srgb is not None else None
            return {"type": "solid", "hex": hex_val}
        if bgPr.find(a_ns + "gradFill") is not None:
            return {"type": "gradient", "hex": None}
        if bgPr.find(a_ns + "pattFill") is not None:
            return {"type": "pattern", "hex": None}
        if bgPr.find(a_ns + "blipFill") is not None:
            return {"type": "blip", "hex": None}
        return {"type": "unknown", "hex": None}

    # 1) slide 自身
    info = parse_bg_from_element(slide._element)
    if info is not None:
        info["source"] = "slide"
        return info
    # 2) 回溯 slideLayout
    try:
        layout = slide.slide_layout
    except Exception:
        layout = None
    info = parse_bg_from_element(layout._element) if layout is not None else None
    if info is not None:
        info["source"] = "layout"
        return info
    # 3) 回溯 slideMaster
    try:
        master = layout.slide_master if layout is not None else None
    except Exception:
        master = None
    info = parse_bg_from_element(master._element) if master is not None else None
    if info is not None:
        info["source"] = "master"
        return info
    return {"type": None, "hex": None, "source": None}


def check_bg_color(prs):
    """+3：第五页和第六页 PPT 背景
        细则要求逐条核验：
          (a) 第 5 页和第 6 页样式与前 4 页一致（比较回溯后的"有效"背景）
          (b) 采用浅蓝色纯色填充
          (c) 颜色代码为 F7FAFC
    """
    n = len(prs.slides)
    infos = [_get_slide_bg_fill_info(prs.slides[i]) for i in range(min(6, n))]

    # 取每页有效背景的 (type, hex) 作为"样式"指纹（忽略来源，允许母版/版式统一设置）
    def style_key(info):
        return (info["type"], (info["hex"] or "").upper() if info["hex"] else None)

    front_keys = [style_key(infos[i]) for i in range(min(4, n))]
    s5_key = style_key(infos[4]) if n >= 5 else None
    s6_key = style_key(infos[5]) if n >= 6 else None

    # (a) 前 4 页样式统一，且 5、6 页有效背景与前 4 页一致
    front_uniform = (
        len(front_keys) >= 1
        and len(set(front_keys)) == 1
        and front_keys[0] != (None, None)
    )
    consistent_with_front_5 = front_uniform and (s5_key == front_keys[0])
    consistent_with_front_6 = front_uniform and (s6_key == front_keys[0])
    style_consistent_ok = consistent_with_front_5 and consistent_with_front_6

    # (b) 第 5、6 页有效背景为纯色填充（回溯后 type=='solid' 也算，
    #     bgRef 引用到主题纯色时也视为纯色）
    def is_solid(info):
        if info is None:
            return False
        if info["type"] == "solid":
            return True
        if info["type"] == "ref" and info["hex"]:
            return True
        return False

    s5_solid_ok = is_solid(infos[4]) if n >= 5 else False
    s6_solid_ok = is_solid(infos[5]) if n >= 6 else False
    solid_fill_ok = s5_solid_ok and s6_solid_ok

    # (c) 颜色代码为 F7FAFC（严格匹配十六进制；纯色或可解析的引用色均可）
    s5_hex_ok = s5_solid_ok and (infos[4]["hex"] or "").upper() == EXPECTED_BG_HEX
    s6_hex_ok = s6_solid_ok and (infos[5]["hex"] or "").upper() == EXPECTED_BG_HEX
    hex_ok = s5_hex_ok and s6_hex_ok

    checks = [
        ("第 5、6 页样式与前 4 页一致", style_consistent_ok),
        ("第 5、6 页采用纯色填充", solid_fill_ok),
        ("颜色代码为 F7FAFC", hex_ok),
    ]
    miss = [name for name, ok in checks if not ok]
    sources = [i["source"] for i in infos]
    detail = (
        f"前4页背景={front_keys}, 第5页={s5_key}, 第6页={s6_key}, 来源={sources}"
    )
    if miss:
        detail += "；不达标：" + "，".join(miss)
        return False, detail
    return True, detail


def check_deductions(prs, file_path):
    """-3：文件中出现批注、红色标记、临时说明、截图边框或多余占位对象。
        细则要求逐类核验，命中任一类即累计扣分依据：
          (a) 批注
          (b) 红色标记
          (c) 临时说明
          (d) 截图边框
          (e) 多余占位对象
    """
    hits = []

    # (a) 批注：检查 ppt/comments*.xml 是否包含真实的 cm 节点
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            for n in zf.namelist():
                base = n.lower().rsplit("/", 1)[-1]
                if base.startswith("comment") and base.endswith(".xml") \
                        and "author" not in base:
                    data = zf.read(n).decode("utf-8", errors="ignore")
                    if "<p:cm " in data or "<p14:cm " in data or "<cm " in data:
                        hits.append(f"(a)批注：发现批注文件 {n}")
                        break
    except Exception:
        pass

    # (b) 红色标记：检查每页中是否出现纯红色文字 / 红色形状描边或填充
    def is_red(hex_str):
        if not hex_str or len(hex_str) < 6:
            return False
        h = hex_str.upper()[-6:]
        try:
            r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        except Exception:
            return False
        # 红色：R 显著大于 G、B，且 R 较亮
        return (r >= 150) and (r - g >= 60) and (r - b >= 60)

    a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            # 文字红色
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    for r in para.runs:
                        try:
                            rgb = r.font.color.rgb
                            if rgb is not None and is_red(str(rgb)):
                                hits.append(f"(b)红色标记：第{i}页文本红色 '{r.text[:20]}'")
                        except Exception:
                            pass
                # 直接扫 XML 中的 srgbClr 兜底
                for srgb in sh.text_frame._txBody.iter(a_ns + "srgbClr"):
                    if is_red(srgb.get("val", "")):
                        hits.append(f"(b)红色标记：第{i}页文本含红色 srgbClr={srgb.get('val')}")
                        break
            # 形状描边/填充红色（仅检查 AUTO_SHAPE / FREEFORM / TEXT_BOX 等非表格非图表对象的 spPr）
            try:
                sp_el = sh._element
                # 仅扫 sh 自身 spPr 的颜色
                for spPr in sp_el.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}spPr"):
                    for srgb in spPr.iter(a_ns + "srgbClr"):
                        if is_red(srgb.get("val", "")):
                            hits.append(f"(b)红色标记：第{i}页形状 {sh.name} 含红色 srgbClr={srgb.get('val')}")
                            break
                    break  # 只看第一个 spPr，避免重复
            except Exception:
                pass

    # (c) 临时说明：搜索可疑文本词
    sus_words = ["临时说明", "TODO", "待修改", "待补充", "占位", "REMOVE", "DELETE", "草稿"]
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if sh.has_text_frame:
                tx = sh.text_frame.text
                for w in sus_words:
                    if w in tx:
                        hits.append(f"(c)临时说明：第{i}页含 '{w}'")

    # (d) 截图边框：检查图片对象的边框描边
    from pptx.shapes.picture import Picture
    for i, sl in enumerate(prs.slides, 1):
        for sh in sl.shapes:
            if isinstance(sh, Picture):
                # 仅当图片显式带有线条/边框时才视为"截图边框"
                has_border = False
                try:
                    sp_el = sh._element
                    for ln in sp_el.iter(a_ns + "ln"):
                        # 有 solidFill / prstDash 等子节点视为有边框
                        if len(list(ln)) > 0 and ln.find(a_ns + "noFill") is None:
                            has_border = True
                            break
                except Exception:
                    pass
                if has_border:
                    hits.append(f"(d)截图边框：第{i}页图片 {sh.name} 含边框")

    # (e) 多余占位对象：第 5、6 页中前 4 页未出现过、且为空内容的对象
    front_names = set()
    for idx in range(min(4, len(prs.slides))):
        for sh in prs.slides[idx].shapes:
            front_names.add(sh.name)
    for i in (5, 6):
        if i - 1 >= len(prs.slides):
            continue
        for sh in prs.slides[i - 1].shapes:
            if sh.has_table:
                continue
            if hasattr(sh, "has_chart") and sh.has_chart:
                continue
            if sh.name in front_names:
                continue
            if sh.has_text_frame and sh.text_frame.text.strip() == "" \
                    and (sh.width or 0) > 0 and (sh.height or 0) > 0:
                hits.append(f"(e)多余占位对象：第{i}页存在非模板的空文本对象 name={sh.name}")

    # 去重
    seen = set()
    uniq = []
    for h in hits:
        if h not in seen:
            seen.add(h)
            uniq.append(h)
    return uniq


# ====== 主流程 ======
def evaluate(dir_path: str) -> dict:
    """统一入口：接收脚本所在目录的路径；脚本自行在该目录内定位并打开被评估文档。

    返回结构见《脚本接口差异与统一建议.md》§2.2。
    """
    # 维度二评分项清单（描述、满分、执行函数）
    dim2_defs = [
        ("PPT页面数量与顺序：最终共6页，第5页为数据表页，第6页为对应折线图页。", 5, check_pages_order),
        ('第5页标题：使用"门店数字化与低碳运营指标"标题，位于页面左上方顶部，字体为微软雅黑24磅加粗。', 1, check_s5_title),
        ("第5页表格结构：包含1行表头和9行数据，共10行8列，所有单元格可编辑。表头从左至右依次为"
         '"年份""智慧门店响应率（%）""线上订单占比（%）""自助结算覆盖率（%）""绿色到店客流占比（%）"'
         '"低碳配送完成率（%）""末端接驳便利度（%）""运营协同指数（0–100）"。'
         "左侧从上到下2016—2018年数据：2016年为31.8、18.9、24.6、27.3、22.6、33.4、29.7；"
         "2017年为34.9、23.7、29.8、31.5、26.1、36.8、34.1；"
         "2018年为42.6、27.9、35.2、34.4、30.6、39.6、39.8。"
         "2019—2021年数据：2019年为45.7、33.8、40.9、37.2、35.9、43.8、44.7；"
         "2020年为53.4、41.7、49.6、42.8、44.3、50.6、52.9；"
         "2021年为58.6、48.8、55.7、47.6、49.7、55.1、58.4。"
         "2022—2024年数据：2022年为66.7、56.9、62.6、55.8、57.4、63.7、66.2；"
         "2023年为71.6、62.4、69.8、61.6、64.9、70.2、72.1；"
         "2024年为78.9、68.1、76.7、67.3、72.6、74.9、78.2。", 5, check_s5_table_structure),
        ("第5页表格样式：使用蓝色或青绿色表头、白色表头文字和浅色交替数据行", 1, check_s5_table_style),
        ("第6页折线图对象：使用PowerPoint原生可编辑折线图，图表数据工作簿中包含第5页全部年份和指标。"
         "横轴从左至右依次显示2016、2017、2018、2019、2020、2021、2022、2023、2024。"
         '数据系列完整包含"智慧门店响应率""线上订单占比""自助结算覆盖率""绿色到店客流占比"'
         '"低碳配送完成率""末端接驳便利度""运营协同指数"7条折线。'
         "7条折线的全部63个数据点与第5页对应单元格完全一致。纵轴采用0至100的统一刻度。"
         '图例准确显示7个系列名称，图例顺序与第5页表头顺序一致，不使用"系列1"等默认名称，'
         "图像数据来源于第五页表格。", 5, check_s6_chart_data),
        ("第6页折线样式：各系列使用可区分的不同颜色线条，并设置清晰的数据点标记。", 1, check_s6_line_style),
        ('第6页图表标题与标注：出现标题且位于页面左上方顶部，'
         '字体为微软雅黑24磅加粗。图表下方有最高点说明内容。', 1, check_s6_title),
        ("第五页和第六页PPT背景：样式和前四页一致，采用浅蓝色纯色填充，颜色代码为F7FAFC", 3, check_bg_color),
    ]
    max_score = sum(pts for _, pts, _ in dim2_defs)

    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": max_score,
    }

    try:
        # 在脚本所在目录内定位被评估文档
        file_path = _locate_pptx(dir_path)
        if not file_path:
            result["status"] = "error"
            result["error"] = f"目录 {dir_path!r} 内未找到 .pptx 文件"
            return result
        result["file_name"] = os.path.basename(file_path)

        prs = Presentation(file_path)
        result["dim1_pass"] = True

        # 维度 2 逐项打分（命中与未命中都列出）
        total = 0
        for desc, pts, fn in dim2_defs:
            try:
                ok, _ = fn(prs)
            except Exception:
                ok = False
            delta = pts if ok else 0
            total += delta
            result["dim2_items"].append({
                "rule": desc,
                "max_delta": pts,
                "delta": delta,
                "hit": bool(ok),
                "detail": "",
            })

        result["total_score"] = total
        return result
    except Exception as e:
        result["status"] = "error"
        result["error"] = f"脚本执行异常：{e}"
        return result


if __name__ == "__main__":
    # 本地调试：默认以脚本所在目录为入参
    _dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(_dir), ensure_ascii=False, indent=2))
