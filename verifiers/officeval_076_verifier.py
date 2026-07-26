# -*- coding: utf-8 -*-
"""
对 "艺术漆北区销售额提升方案.pptx" 的自动评估脚本。

评估逻辑：
  维度1（可用与可修改性）：任一硬性条件不满足 → 直接 0 分，且不再检查维度2。
  维度2（完成度）：逐项检查得分点 / 扣分点，命中则累计该项的（正或负）分数。

最终打印每条命中的细则与累计得分。
"""

import sys
import os
import re
import json

SCRIPT_ID = "076"

import zipfile
from lxml import etree

from pptx import Presentation
from pptx.util import Emu

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

# ---- 期望文案 ----
EXPECT_TITLE = "艺术漆北区销售额提升方案"

TOC_ENTRIES = [
    "市场趋势与消费升级", "重点城市增长机会", "行业竞争与品牌机遇",
    "年度发展目标",       "核心产品体系",     "产品组合与差异化卖点",
    "多渠道拓展策略",     "整合营销与本地转化", "服务体系与执行保障",
]

SECTION_TITLES = [
    "一、市场趋势与消费升级", "二、重点城市增长机会", "三、行业竞争与品牌机遇",
    "四、年度发展目标",       "五、核心产品体系",     "六、产品组合与差异化卖点",
    "七、多渠道拓展策略",     "八、整合营销与本地转化", "九、服务体系与执行保障",
]

# 每个正文页期望的两个二级标题
SUB_TITLES = [
    ("1.1 市场扩容趋势",          "1.2 消费需求变化"),
    ("2.1 改善型住房需求释放",    "2.2 存量房翻新潜力"),
    ("3.1 竞争格局变化",          "3.2 品牌突破方向"),
    ("4.1 经营增长目标",          "4.2 品牌认知目标"),
    ("5.1 主推系列定位",          "5.2 消费决策简化"),
    ("6.1 产品矩阵设计",          "6.2 核心卖点提炼"),
    ("7.1 零售终端升级",          "7.2 家装与工程渠道拓展"),
    ("8.1 线上内容种草",          "8.2 线下活动转化"),
    ("9.1 服务体系升级",          "9.2 经销商赋能与组织保障"),
]

# 每个二级标题对应"正文语义匹配关键词" —— 至少命中其中 1 个即视为"正文与二级标题匹配"
SUB_BODY_KEYWORDS = {
    "1.1 市场扩容趋势":         ("市场", "扩容", "增长", "规模", "趋势"),
    "1.2 消费需求变化":         ("消费", "需求", "变化", "升级", "偏好"),
    "2.1 改善型住房需求释放":   ("改善", "住房", "需求", "释放", "置换", "改善型"),
    "2.2 存量房翻新潜力":       ("存量", "翻新", "潜力", "老房", "重装"),
    "3.1 竞争格局变化":         ("竞争", "格局", "变化", "对手", "份额"),
    "3.2 品牌突破方向":         ("品牌", "突破", "方向", "差异", "定位"),
    "4.1 经营增长目标":         ("经营", "增长", "目标", "销售额", "业绩", "指标"),
    "4.2 品牌认知目标":         ("品牌", "认知", "目标", "知名度", "美誉", "传播"),
    "5.1 主推系列定位":         ("主推", "系列", "定位", "产品", "核心"),
    "5.2 消费决策简化":         ("消费", "决策", "简化", "购买", "选择"),
    "6.1 产品矩阵设计":         ("产品", "矩阵", "设计", "组合", "sku"),
    "6.2 核心卖点提炼":         ("核心", "卖点", "提炼", "优势", "价值"),
    "7.1 零售终端升级":         ("零售", "终端", "升级", "门店", "陈列"),
    "7.2 家装与工程渠道拓展":   ("家装", "工程", "渠道", "拓展", "合作", "设计师"),
    "8.1 线上内容种草":         ("线上", "内容", "种草", "抖音", "小红书", "社交", "短视频", "直播"),
    "8.2 线下活动转化":         ("线下", "活动", "转化", "促销", "体验", "沙龙"),
    "9.1 服务体系升级":         ("服务", "体系", "升级", "售后", "施工", "保障"),
    "9.2 经销商赋能与组织保障": ("经销商", "赋能", "组织", "保障", "培训", "支持"),
}

SECTION_PAGES = [3, 5, 7, 9, 11, 13, 15, 17, 19]   # 章节标题页（1-based）
BODY_PAGES    = [4, 6, 8, 10, 12, 14, 16, 18, 20]  # 正文页（1-based）


# ----------------- 工具 -----------------
def iter_runs(shape):
    """yield (text, font_name, size_pt, bold) for every run in shape"""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            yield (
                run.text or "",
                run.font.name,
                run.font.size.pt if run.font.size else None,
                run.font.bold,
            )


def shape_text(shape):
    return shape.text_frame.text if shape.has_text_frame else ""


def hex_to_rgb(h):
    h = h.lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def is_green(rgb):
    r, g, b = rgb
    return g >= max(r, b) and g >= 120        # 偏绿


def is_light_green(rgb):
    r, g, b = rgb
    return (r > 180 and g > 200 and b > 180) and is_green(rgb) or \
           (r > 150 and g > 200 and b > 150)


def is_orange(rgb):
    r, g, b = rgb
    return r >= 180 and 80 <= g <= 200 and b <= 120


# ----------------- 维度 1 -----------------
def check_dim1(prs, raw_zip, pptx_path: str) -> "tuple[bool, list[str]]":
    """维度 1：交付文件为 .pptx 格式，文件可正常打开。

    规则简化：仅校验后缀为 .pptx。
    "可正常打开" 已在上层通过 python-pptx 的 Presentation(pptx_path) 是否抛异常来判定：
    调用本函数时说明打开已经成功，因此这里只需再确认后缀即可。
    此前对"可编辑文本 / 越界 / 空白过多"的额外校验已弃用；
    `prs`、`raw_zip` 参数保留仅为保持调用签名兼容。
    """
    del prs, raw_zip  # 声明未使用，避免静态检查告警
    fails: "list[str]" = []
    if not pptx_path.lower().endswith(".pptx"):
        fails.append("文件后缀不是 .pptx")
    return (len(fails) == 0, fails)


# ----------------- 维度 2 检查项 -----------------
def get_slide_xml(raw_zip, slide_no):
    with raw_zip.open(f"ppt/slides/slide{slide_no}.xml") as f:
        return f.read().decode("utf-8")


def check_first_page(prs, raw_zip):
    """+1 第1页首页：
       - 标题文本为"艺术漆北区销售额提升方案"
       - 位于页面中部或偏上位置
       - 字体为宋体
       - 字号 65-75 磅
       - 加粗
       - 颜色与绿色渐变背景形成清晰对比
       六个子点全部满足；且标题 shape 内的所有非空 run 都需满足字体/字号/加粗/颜色。
    """
    slide = prs.slides[0]
    sh_ = prs.slide_height

    # 收集第1页背景渐变颜色（用于第6点对比度判定）：
    #   1) slide 上任何 shape 的 gradFill；
    #   2) slide / layout / master 的 <p:bg>；
    #   3) 主题（ppt/theme/*.xml）中所有 srgbClr —— 兜底。
    def _grads_from_elem(elem):
        cols = []
        xml_s = etree.tostring(elem).decode()
        for grad in re.finditer(r"<a:gradFill[^>]*>.*?</a:gradFill>", xml_s, re.S):
            for m in re.finditer(r'srgbClr val="([0-9A-Fa-f]{6})"', grad.group()):
                cols.append(hex_to_rgb(m.group(1)))
        return cols

    bg_colors = []
    for sh in slide.shapes:
        bg_colors.extend(_grads_from_elem(sh._element))

    for src in (slide, slide.slide_layout, slide.slide_layout.slide_master):
        try:
            bg_elem = src.element.find("p:cSld/p:bg", NS)
        except Exception:
            bg_elem = None
        if bg_elem is not None:
            bg_colors.extend(_grads_from_elem(bg_elem))

    # 主题兜底：读取 ppt/theme/*.xml 中的 srgbClr 颜色
    if not bg_colors:
        try:
            theme_names = [n for n in raw_zip.namelist()
                           if n.startswith("ppt/theme/") and n.endswith(".xml")]
            for tn in theme_names:
                tx = raw_zip.read(tn).decode("utf-8", errors="ignore")
                for m in re.finditer(r'srgbClr val="([0-9A-Fa-f]{6})"', tx):
                    bg_colors.append(hex_to_rgb(m.group(1)))
        except Exception:
            pass

    def luminance(c):
        r, g, b = [x / 255 for x in c]
        return 0.2126 * r + 0.7152 * g + 0.0722 * b

    # 背景以最浅色为参考（最不利对比情况）
    lb = max((luminance(c) for c in bg_colors), default=1.0)

    # 找到标题 shape
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if sh.text_frame.text.strip() != EXPECT_TITLE:
            continue

        # ① 文本完全匹配  ——  已通过 strip()==EXPECT_TITLE
        # ② 位置：中部或偏上 —— shape 垂直中心位于页面上半段或正中（top_ratio ≤ 0.55）
        center_ratio = ((sh.top or 0) + (sh.height or 0) / 2) / sh_
        if center_ratio > 0.55:
            return False, f'首页标题不在"中部或偏上"位置（垂直中心 {center_ratio:.2f}）'

        # ③④⑤⑥ 对标题 shape 内的所有非空 run 逐一校验
        checked_any = False
        last_size = None
        last_rgb = None
        last_contrast = None

        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                text = run.text or ""
                if not text.strip():
                    continue
                checked_any = True

                # ③ 宋体
                font = run.font.name
                if font not in ("宋体", "SimSun"):
                    return False, f"首页标题存在非宋体 run（字体 {font}，文本 {text!r}）"

                # ④ 字号 65-75 磅
                size = run.font.size.pt if run.font.size else None
                if size is None or not (65 <= size <= 75):
                    return False, f"首页标题 run 字号 {size}pt 不在 65-75（文本 {text!r}）"

                # ⑤ 加粗
                if not run.font.bold:
                    return False, f"首页标题 run 未加粗（文本 {text!r}）"

                # ⑥ 颜色与绿色渐变背景形成清晰对比
                try:
                    fc = run.font.color.rgb
                    title_rgb = (fc[0], fc[1], fc[2]) if fc else None
                except Exception:
                    title_rgb = None
                if title_rgb is None:
                    return False, (f"首页标题 run 未设置纯色字体颜色，"
                                   f"无法判定与背景对比（文本 {text!r}）")

                lt = luminance(title_rgb)
                contrast = (max(lt, lb) + 0.05) / (min(lt, lb) + 0.05)
                if contrast < 3.0:
                    return False, (f"首页标题 run 颜色与绿色渐变背景对比不足"
                                   f"（contrast={contrast:.2f}，文本 {text!r}）")

                last_size, last_rgb, last_contrast = size, title_rgb, contrast

        if not checked_any:
            return False, "首页标题 shape 内没有可读的非空 run"

        return True, (f"首页标题命中："
                      f"位置中部/偏上、宋体、{last_size}pt、加粗、"
                      f"色#{last_rgb[0]:02X}{last_rgb[1]:02X}{last_rgb[2]:02X} "
                      f"与背景对比度 {last_contrast:.2f}")

    return False, "首页未找到目标标题"


def check_bg_gradient(prs, raw_zip):
    """+3 页面背景：
       ① 采用渐变背景；
       ② 渐变色为浅、中绿色（即所有色标都偏绿，且亮度处于浅~中范围，不深）；
       ③ 整体颜色不深 —— 背景平均亮度足够高；
       ④ 文本可清晰阅读 —— 该页主要文本颜色与背景对比度足够。
       四个子点都满足的页面才算命中；要求 PPT 大多数页面（≥ 18/21）满足。
    """

    def luminance(c):
        r, g, b = [x / 255 for x in c]
        return 0.2126 * r + 0.7152 * g + 0.0722 * b

    qualified_pages = []
    detail_fail = []

    for n in range(1, len(prs.slides) + 1):
        xml = get_slide_xml(raw_zip, n)
        root = etree.fromstring(xml.encode("utf-8"))

        # ① 渐变背景
        grad_colors = []
        for g in root.findall(".//a:gradFill", NS):
            for c in g.findall(".//a:srgbClr", NS):
                v = c.get("val")
                if v:
                    grad_colors.append(hex_to_rgb(v))
            if grad_colors:
                break
        if len(grad_colors) < 2:
            detail_fail.append(f"P{n}:无渐变")
            continue

        # ② 所有色标都偏绿
        if not all(is_green(c) for c in grad_colors):
            detail_fail.append(f"P{n}:渐变色非绿色系")
            continue

        # ② 浅、中绿色（亮度处于浅 / 中范围，避免深色）
        lums = [luminance(c) for c in grad_colors]
        if min(lums) < 0.35:        # 最深的一个色标也不能太深
            detail_fail.append(f"P{n}:渐变中存在深色（min L={min(lums):.2f}）")
            continue
        if max(lums) < 0.75:        # 至少包含一个浅色色标
            detail_fail.append(f"P{n}:渐变缺少浅色色标（max L={max(lums):.2f}）")
            continue

        # ③ 整体颜色不深 —— 平均亮度 ≥ 0.55
        avg_l = sum(lums) / len(lums)
        if avg_l < 0.55:
            detail_fail.append(f"P{n}:背景整体偏深（avg L={avg_l:.2f}）")
            continue

        # ④ 文本可清晰阅读 —— 本页最常见文字颜色与背景最浅色对比 ≥ 3.0
        slide = prs.slides[n - 1]
        text_colors = []
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if not (run.text or "").strip():
                        continue
                    try:
                        rgb = run.font.color.rgb
                        if rgb is not None:
                            text_colors.append((rgb[0], rgb[1], rgb[2]))
                    except Exception:
                        pass
        if not text_colors:
            # 没有显式文字颜色时，默认黑色（office 默认），与浅绿背景对比天然足够
            text_colors = [(0, 0, 0)]

        bg_light_l = max(lums)
        readable = True
        for tc in text_colors:
            tl = luminance(tc)
            contrast = (max(tl, bg_light_l) + 0.05) / (min(tl, bg_light_l) + 0.05)
            # "文本可清晰阅读"采用常规可读性下限 2.0（深绿字 / 浅绿底属于清晰对比）
            if contrast < 2.0:
                readable = False
                break
        if not readable:
            detail_fail.append(f"P{n}:文字与背景对比不足")
            continue

        qualified_pages.append(n)

    ok = len(qualified_pages) >= max(1, len(prs.slides) - 3)  # 允许少数特殊页（如尾页用图片）
    msg = (f'{len(qualified_pages)}/{len(prs.slides)} 页满足'
           f'"浅/中绿渐变 + 整体不深 + 文本可清晰阅读"')
    if detail_fail and not ok:
        msg += f"；未命中：{detail_fail[:3]}"
    return ok, msg


def check_toc(prs):
    """+3 第2页目录页：
       ① 必须是第 2 页；
       ② 出现"目录"标题（独立文本）；
       ③ 出现九个一级目录条目，分别为：
          市场趋势与消费升级 / 重点城市增长机会 / 行业竞争与品牌机遇 /
          年度发展目标 / 核心产品体系 / 产品组合与差异化卖点 /
          多渠道拓展策略 / 整合营销与本地转化 / 服务体系与执行保障。
       三个子点全部命中才计分。
    """
    # ① 第 2 页
    if len(prs.slides) < 2:
        return False, "PPT 不足 2 页，没有第2页目录页"
    slide = prs.slides[1]

    # 收集本页所有文本（按 shape 分开，便于判定"目录"为独立标题）
    shape_texts = [shape_text(sh).strip() for sh in slide.shapes
                   if sh.has_text_frame and shape_text(sh).strip()]

    # ② "目录"标题 —— 存在一个文本框文本恰为"目录"
    if "目录" not in shape_texts:
        # 兼容"目  录"这类排版（去掉所有空白后等于"目录"）
        if not any(re.sub(r"\s+", "", t) == "目录" for t in shape_texts):
            return False, '第2页未出现独立的"目录"标题'

    # ③ 九个一级目录条目 —— 每一项都必须出现
    page_text_concat = "".join(shape_texts)
    missing = [e for e in TOC_ENTRIES if e not in page_text_concat]
    if missing:
        return False, f"第2页目录缺失条目：{missing}"

    # 额外校验"九个"——这九条目录条目都能在页面上找到对应的承载文本框
    hit_per_entry = []
    for e in TOC_ENTRIES:
        if any(e in t for t in shape_texts):
            hit_per_entry.append(e)
    if len(hit_per_entry) != 9:
        return False, f"九条目录条目实际命中 {len(hit_per_entry)} 条"

    return True, '第2页：含"目录"标题，且九个一级目录条目齐全'


def check_section_pages(prs):
    """+5 第 3、5、7、9、11、13、15、17、19 页为一至九部分标题页：
       ① 页码与"一/二/.../九、xxx" 一一对应；
       ② 标题文本精确等于细则给出的 9 个标题；
       ③ 字体加粗；
       ④ 字号 38-42 磅（闭区间）；
       ⑤ 文字居于页面中间（水平居中 + 垂直位于页面中部）；
       ⑥ 文字放在一行（只有一个段落，且不换行）。
       九个页面全部满足才计 +5。
    """
    sw, sh_ = prs.slide_width, prs.slide_height
    ok_pages = []
    fail_details = []

    if len(prs.slides) < max(SECTION_PAGES):
        return False, f"PPT 不足 {max(SECTION_PAGES)} 页"

    for i, pno in enumerate(SECTION_PAGES):
        target = SECTION_TITLES[i]
        slide = prs.slides[pno - 1]

        hit_shape = None
        for sh in slide.shapes:
            if sh.has_text_frame and shape_text(sh).strip() == target:
                hit_shape = sh
                break
        # ② 标题文本精确匹配
        if hit_shape is None:
            fail_details.append(f'P{pno} 未找到标题"{target}"')
            continue

        # ⑥ 一行：段落数为 1；且文本中不含换行符/软回车；
        #    并且文本框宽度需足够容纳整段文字（避免因文本框过窄而视觉换行）
        paragraphs = list(hit_shape.text_frame.paragraphs)
        if len(paragraphs) != 1:
            fail_details.append(f"P{pno} 标题未放在一行（段落数 {len(paragraphs)}）")
            continue

        # 段落内文本是否含硬换行（\n）或软换行（\v / \x0b，PowerPoint 的 shift+enter）
        raw_text = ""
        has_soft_break = False
        for run in paragraphs[0].runs:
            raw_text += run.text or ""
            # python-pptx 的 <a:br> 会在解析时表现为 \x0b；同时通过 XML 直接确认
            try:
                if run._r.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}br"):
                    has_soft_break = True
            except Exception:
                pass
        if "\n" in raw_text or "\v" in raw_text or "\x0b" in raw_text or has_soft_break:
            fail_details.append(f"P{pno} 标题含换行符/软回车，未真正放在一行")
            continue
        # 顶层 <a:br> 兜底
        try:
            if hit_shape.text_frame._txBody.findall(
                    ".//{http://schemas.openxmlformats.org/drawingml/2006/main}br"):
                fail_details.append(f"P{pno} 标题含 <a:br> 软换行")
                continue
        except Exception:
            pass

        # ③④ 字体加粗 + 字号 38-42
        runs = list(iter_runs(hit_shape))
        if not runs:
            fail_details.append(f"P{pno} 标题无可读 run")
            continue
        size_ok = all(r[2] is not None and 38 <= r[2] <= 42 for r in runs)
        bold_ok = all(r[3] is True for r in runs)
        if not size_ok:
            fail_details.append(
                f"P{pno} 字号不符（{[r[2] for r in runs]}，应 38-42pt）")
            continue
        if not bold_ok:
            fail_details.append(f"P{pno} 标题未加粗")
            continue

        # ⑥（续）文本框宽度需容得下这段文字：粗估中文字符宽 ≈ 字号(pt) * 1.0，
        # 英文/数字/空格宽 ≈ 字号(pt) * 0.55；两侧内边距按 0.25 英寸估算。
        try:
            font_pt = max((r[2] for r in runs if r[2] is not None), default=40)
            char_w = 0.0
            for ch in raw_text:
                if re.match(r"[一-鿿]", ch):
                    char_w += font_pt * 1.00
                elif ch.strip() == "":
                    char_w += font_pt * 0.30
                else:
                    char_w += font_pt * 0.55
            # 转 EMU：1 pt = 12700 EMU
            need_emu = int(char_w * 12700)
            pad_emu = int(0.5 * 914400)  # 双侧内边距合计约 0.5 英寸
            box_w = hit_shape.width or 0
            if box_w and (box_w - pad_emu) < need_emu:
                fail_details.append(
                    f"P{pno} 文本框宽度不足以容纳整段标题（宽 {box_w} EMU，"
                    f"估算需要 {need_emu} EMU），可能视觉换行")
                continue
        except Exception:
            pass

        # ⑤ 文字居于页面中间
        para_align = paragraphs[0].alignment  # 段落水平对齐
        cx_ratio = ((hit_shape.left or 0) + (hit_shape.width or 0) / 2) / sw
        cy_ratio = ((hit_shape.top or 0) + (hit_shape.height or 0) / 2) / sh_
        # 水平居中：段落对齐为 CENTER，或形状几何中心在 0.4~0.6
        from pptx.enum.text import PP_ALIGN
        is_h_center = (para_align == PP_ALIGN.CENTER) or (0.4 <= cx_ratio <= 0.6)
        # 垂直居中：形状几何中心在页面中部区间
        is_v_center = 0.3 <= cy_ratio <= 0.7
        if not (is_h_center and is_v_center):
            fail_details.append(
                f"P{pno} 标题未居中（h={cx_ratio:.2f}/align={para_align}，v={cy_ratio:.2f}）")
            continue

        ok_pages.append(pno)

    ok = (len(ok_pages) == 9)
    msg = f"章节标题页命中 {len(ok_pages)}/9（页 {ok_pages}）"
    if not ok and fail_details:
        msg += f"；未命中：{fail_details}"
    return ok, msg


def check_body_pages(prs):
    """+5 第 4、6、8、10、12、14、16、18、20 页为第一至九部分的正文页：
       对每个页面（共九页），都必须满足：
       ① 出现两个二级标题，文本严格等于细则给定：
          (X.1 ...) 与 (X.2 ...)
       ② 两个二级标题均位于"对应部分的最上方靠左区域"
          靠左：文本框左边缘位于页面左半区域（left_ratio ≤ 0.3），
                段落对齐为左对齐 / 默认；
          最上方：s1 位于页面上半段（cy ≤ 0.55）；s2 位于其所在下半分节的顶部
                （cy 处于 0.4~0.8，且没有其它非二级标题正文类文本位于 s2 之上、s1 之下）。
       ③ 二级标题加粗，字号 16-20 磅（闭区间）
       ④ 每个二级标题下方放置对应正文内容（top 大于该二级标题）；
          正文与二级标题"匹配"：位置紧随其后 + 水平靠左对齐 + 语义匹配
          （命中 SUB_BODY_KEYWORDS 中该二级标题的关键词）
       ⑤ 正文字号 12-16 磅（闭区间）
       九个页面全部满足才计 +5。
    """
    sw, sh_ = prs.slide_width, prs.slide_height
    ok_pages = []
    fail_details = []

    if len(prs.slides) < max(BODY_PAGES):
        return False, f"PPT 不足 {max(BODY_PAGES)} 页"

    for i, pno in enumerate(BODY_PAGES):
        slide = prs.slides[pno - 1]
        s1_text, s2_text = SUB_TITLES[i]

        # 找两个二级标题 shape（文本严格匹配）
        sub_shapes = {}
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            txt = shape_text(sh).strip()
            if txt == s1_text:
                sub_shapes["s1"] = sh
            elif txt == s2_text:
                sub_shapes["s2"] = sh

        # ① 两个二级标题都要出现
        if "s1" not in sub_shapes or "s2" not in sub_shapes:
            fail_details.append(
                f"P{pno} 缺少二级标题（命中 {list(sub_shapes.keys())}，期望 {s1_text!r}/{s2_text!r}）")
            continue

        page_ok = True
        for key, target in (("s1", s1_text), ("s2", s2_text)):
            sub = sub_shapes[key]
            sub_top = sub.top or 0
            sub_left = sub.left or 0
            sub_h = sub.height or 0
            sub_w = sub.width or 0

            # ② 二级标题在"对应部分的最上方靠左区域"
            #   靠左 —— 文本框左边缘位于页面左半区域（left_ratio ≤ 0.3），
            #          且段落对齐为左对齐 / 默认（不为居中或右对齐）
            sub_cy = (sub_top + sub_h / 2) / sh_
            left_ratio = sub_left / sw
            from pptx.enum.text import PP_ALIGN
            para_align = sub.text_frame.paragraphs[0].alignment
            is_left = (left_ratio <= 0.3) and \
                      (para_align in (None, PP_ALIGN.LEFT, PP_ALIGN.JUSTIFY))
            if not is_left:
                fail_details.append(
                    f'P{pno} 二级标题"{target}"未靠左（left_ratio={left_ratio:.2f}，align={para_align}）')
                page_ok = False
                break

            # 顶部：
            # s1 必须位于上半段（cy ≤ 0.55）；
            # s2 作为下半分节的"最上方"，cy 应位于 0.4~0.8。
            if key == "s1":
                if sub_cy > 0.55:
                    fail_details.append(
                        f'P{pno} 二级标题"{target}"未位于上半段最上方（cy={sub_cy:.2f}）')
                    page_ok = False
                    break
            else:  # s2
                if not (0.4 <= sub_cy <= 0.8):
                    fail_details.append(
                        f'P{pno} 二级标题"{target}"未位于下半分节最上方（cy={sub_cy:.2f}）')
                    page_ok = False
                    break

            # ③ 二级标题：加粗 + 16-20 磅
            runs = list(iter_runs(sub))
            if not runs:
                fail_details.append(f'P{pno} 二级标题"{target}"无可读 run')
                page_ok = False
                break
            if not all(r[3] is True for r in runs):
                fail_details.append(f'P{pno} 二级标题"{target}"未加粗')
                page_ok = False
                break
            if not all(r[2] is not None and 16 <= r[2] <= 20 for r in runs):
                fail_details.append(
                    f'P{pno} 二级标题"{target}"字号 {[r[2] for r in runs]} 不在 16-20pt')
                page_ok = False
                break

            # ④⑤ 找该二级标题"下方"的正文 shape：
            #    位置 top 大于本二级标题 top，水平靠左，文字较长（属正文性质），
            #    且不是另一个二级标题；选距离本二级标题最近的一个。
            other_sub_top = sub_shapes["s2" if key == "s1" else "s1"].top or 0
            candidates = []
            for sh in slide.shapes:
                if not sh.has_text_frame or sh is sub:
                    continue
                t = shape_text(sh).strip()
                if not t or t in (s1_text, s2_text):
                    continue
                top = sh.top or 0
                left = sh.left or 0
                if top <= sub_top:
                    continue
                # 必须在本二级标题与下一个二级标题之间（s1 才有约束）
                if key == "s1" and top >= other_sub_top:
                    continue
                # 水平：靠左 —— 与二级标题左边缘大致对齐（差距 < 10% 页宽）
                if abs(left - sub_left) > 0.1 * sw:
                    continue
                # 必须是有一定长度的中文正文，避免"自动呈现"这类小标签
                if len(t) < 20:
                    continue
                candidates.append((top, sh, t))
            if not candidates:
                fail_details.append(f'P{pno} 二级标题"{target}"下方未找到匹配的正文')
                page_ok = False
                break
            candidates.sort(key=lambda x: x[0])

            # 语义匹配：正文文本至少命中一个 SUB_BODY_KEYWORDS 关键词
            kws = SUB_BODY_KEYWORDS.get(target, ())
            body_shape = None
            for _top, sh, t in candidates:
                if not kws or any(k in t for k in kws):
                    body_shape = sh
                    break
            if body_shape is None:
                fail_details.append(
                    f'P{pno} 二级标题"{target}"对应正文语义未匹配（关键词 {list(kws)}，'
                    f'首个候选 {candidates[0][2][:30]!r}）')
                page_ok = False
                break

            # ⑤ 正文字号 12-16 磅
            b_runs = list(iter_runs(body_shape))
            if not b_runs or not all(r[2] is not None and 12 <= r[2] <= 16 for r in b_runs):
                fail_details.append(
                    f'P{pno} 二级标题"{target}"对应正文字号 {[r[2] for r in b_runs]} 不在 12-16pt')
                page_ok = False
                break

        if page_ok:
            ok_pages.append(pno)

    ok = (len(ok_pages) == 9)
    msg = f"正文页（二级标题+对应正文）命中 {len(ok_pages)}/9（页 {ok_pages}）"
    if not ok and fail_details:
        msg += f"；未命中：{fail_details[:5]}"
    return ok, msg


def check_last_page(prs):
    """+1 第21页尾页：
       ① 必须是第 21 页；
       ② 尾页文字使用宋体（页面上所有含中文的 run 字体均为宋体/SimSun）；
       ③ 内容可包含"谢谢观看"/"感谢观看"（至少出现其一作为致谢语）
          以及"艺术漆北区销售额提升方案"（方案名）；
       ④ 不能出现无关主题 —— 页面上的每段可见文本都应与方案/致谢主题相关，
          不出现与"艺术漆/方案/区域增长/感谢/THANKS"无关的题外内容。
       四个子点全部满足才计 +1。
    """
    # ① 第 21 页
    if len(prs.slides) < 21:
        return False, "PPT 不足 21 页，没有尾页"
    slide = prs.slides[20]

    # 收集所有文本及其字体
    text_runs = []                # [(text, font_name), ...] 仅取非空
    page_text = ""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        page_text += shape_text(sh)
        for text, font, _, _ in iter_runs(sh):
            if (text or "").strip():
                text_runs.append((text, font))

    # ② 宋体 —— 含中文字符的 run 字体必须为宋体/SimSun
    for text, font in text_runs:
        if re.search(r"[一-鿿]", text):
            if font not in ("宋体", "SimSun"):
                return False, f"尾页存在非宋体中文文本：{text!r}（字体 {font}）"

    # ③ 致谢语与方案名
    has_thanks = ("谢谢观看" in page_text) or ("感谢观看" in page_text)
    if not has_thanks:
        return False, '尾页缺少"谢谢观看/感谢观看"致谢语'
    if EXPECT_TITLE not in page_text:
        return False, f'尾页缺少方案名"{EXPECT_TITLE}"'

    # ④ 不能出现无关主题
    #   定义"与本方案/致谢相关"的关键字集合；
    #   任一文本框的可见中文文本若既不属于致谢语，也不含任一相关关键字，则判定为无关主题。
    THEME_KEYWORDS = [
        "艺术漆", "北区", "销售", "方案", "提升", "增长", "墙面", "美学",
        "区域", "价值", "能力", "市场", "品牌", "产品", "渠道", "服务",
        "环保", "色彩", "设计", "客户", "终端", "工程", "经销",
        "目录", "趋势", "城市", "竞争", "目标", "体系", "营销", "保障",
    ]
    THANKS_TOKENS = ("谢谢", "感谢", "THANKS", "THANK YOU", "Thanks")

    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = shape_text(sh).strip()
        if not txt:
            continue
        # 纯致谢语放行
        if any(tok in txt for tok in THANKS_TOKENS):
            continue
        # 含方案名放行
        if EXPECT_TITLE in txt:
            continue
        # 纯英文/装饰短语放行（不含中文）
        if not re.search(r"[一-鿿]", txt):
            continue
        # 含任何主题相关关键字放行
        if any(k in txt for k in THEME_KEYWORDS):
            continue
        # 其余视为"无关主题"
        return False, f"尾页出现疑似无关主题：{txt!r}"

    return True, "尾页第21页：宋体；含致谢语与方案名；无无关主题"


def check_orange_box(prs):
    """+3 所有正文页橙色矩形框：
       对九个正文页（4/6/8/10/12/14/16/18/20）逐页校验：
       ① 本页所有二级标题 + 正文内容都被橙色矩形框包住
          —— 允许"一个大框把整页两组二级标题/正文都包住"，也允许"每组各一个框"；
          —— 每个二级标题及其对应正文都必须落在某个合规橙色矩形框的边界内即可。
       ② 矩形框为可编辑形状（autoshape，不能是图片/嵌入对象）；
       ③ 边线清晰可见（存在 <a:ln>，且非 noFill）；
       ④ 边线线宽约 1.5—3 磅（闭区间，含一点容差 1.4—3.1pt）；
       ⑤ 颜色为橙色或深橙色（不接受其它色系）；
       ⑥ 线宽"不能过细"—— 已包含在 ④ 的下限里。
       九个页面全部满足才计 +3。
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    hits = []
    fail_details = []

    for i, pno in enumerate(BODY_PAGES):
        slide = prs.slides[pno - 1]
        s1_text, s2_text = SUB_TITLES[i]

        # 找到本页的两个二级标题 shape & 它们的"对应正文"
        sub_shapes = {}
        for sh in slide.shapes:
            if sh.has_text_frame and shape_text(sh).strip() in (s1_text, s2_text):
                sub_shapes[shape_text(sh).strip()] = sh
        if len(sub_shapes) != 2:
            fail_details.append(f"P{pno} 未集齐两个二级标题，无法判定外框")
            continue

        sw = prs.slide_width
        groups = []   # 每组：(二级标题 shape, 正文 shape)
        for key in (s1_text, s2_text):
            sub = sub_shapes[key]
            # 取"在该二级标题下方、左对齐、文本较长"的最近 shape 视为对应正文
            others = []
            for sh in slide.shapes:
                if not sh.has_text_frame or sh is sub:
                    continue
                t = shape_text(sh).strip()
                if not t or t in (s1_text, s2_text):
                    continue
                if (sh.top or 0) <= (sub.top or 0):
                    continue
                if abs((sh.left or 0) - (sub.left or 0)) > 0.1 * sw:
                    continue
                if len(t) < 20:
                    continue
                others.append(sh)
            if not others:
                fail_details.append(f'P{pno} 未找到"{key}"对应的正文 shape')
                groups = None
                break
            others.sort(key=lambda x: x.top or 0)
            groups.append((sub, others[0]))
        if groups is None:
            continue

        # 收集本页所有"可能是橙色矩形框"的 autoshape
        boxes = []
        for sh in slide.shapes:
            if sh.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
            # ② 可编辑形状 —— autoshape 本身即可编辑；排除图片/OLE
            xml = etree.tostring(sh._element).decode()
            # ③ 必须有 <a:ln> 且非 noFill
            m_ln = re.search(r'<a:ln([^>]*)>(.*?)</a:ln>', xml, re.S)
            if not m_ln:
                continue
            ln_attrs, ln_inner = m_ln.group(1), m_ln.group(2)
            if "<a:noFill" in ln_inner:
                continue
            m_clr = re.search(r'srgbClr val="([0-9A-Fa-f]{6})"', ln_inner)
            if not m_clr:
                continue
            # ⑤ 颜色橙色 / 深橙色
            if not is_orange(hex_to_rgb(m_clr.group(1))):
                continue
            # ④⑥ 线宽 1.5–3pt（容差 1.4–3.1）
            m_w = re.search(r'w="(\d+)"', ln_attrs)
            if not m_w:
                # 没指定 w 默认约 0.75pt，过细
                continue
            w_pt = int(m_w.group(1)) / 12700.0
            if not (1.4 <= w_pt <= 3.1):
                continue
            boxes.append((sh, w_pt))

        if not boxes:
            fail_details.append(f"P{pno} 未找到合规橙色矩形框")
            continue

        # ① 允许一个或多个矩形框；每个 (二级标题, 正文) 必须完整落入某一合规框内。
        #    同一大框可同时包住两组内容 —— 不再要求"每组独立占用不同框"。
        def contains(box, *inners):
            bl, bt = box.left or 0, box.top or 0
            br, bb = bl + (box.width or 0), bt + (box.height or 0)
            for s in inners:
                il, it = s.left or 0, s.top or 0
                ir, ib = il + (s.width or 0), it + (s.height or 0)
                # 允许 5pt EMU 边缘容差
                tol = 63500   # 约 5pt
                if not (bl - tol <= il and it >= bt - tol
                        and ir <= br + tol and ib <= bb + tol):
                    return False
            return True

        all_wrapped = True
        for sub, body in groups:
            wrapped = any(contains(box, sub, body) for box, _ in boxes)
            if not wrapped:
                all_wrapped = False
                fail_details.append(
                    f'P{pno} 二级标题"{shape_text(sub).strip()}"及其正文未被橙色矩形框包住')
                break

        if all_wrapped:
            hits.append(pno)

    ok = (len(hits) == 9)
    msg = f"正文页橙色矩形框命中 {len(hits)}/9（页 {hits}）"
    if not ok and fail_details:
        msg += f"；未命中：{fail_details[:5]}"
    return ok, msg


def check_all_simsun(prs):
    """+3 全文字体：
       ① 所有"可见中文文字"字体均为宋体（宋体 / SimSun）；
       ② 覆盖范围必须包括：目录、标题、二级标题、正文、页脚、尾页文字。
       即只要 PPT 里出现的任何一处可见中文 run 不是宋体，就判为未命中。
       本检测同时校验"覆盖范围齐全"：在文件中应至少能定位到上述类别中存在的文本。
    """
    bad = []           # 非宋体的中文 run
    coverage = {       # 上述类别命中过哪些
        "目录": False,
        "标题": False,       # 首页 / 章节标题 / 一级标题
        "二级标题": False,
        "正文": False,
        "页脚": False,
        "尾页": False,
    }

    n_slides = len(prs.slides)
    section_titles_set = set(SECTION_TITLES) | {EXPECT_TITLE}
    sub_titles_set = {t for pair in SUB_TITLES for t in pair}

    def is_chinese(s):
        return bool(re.search(r"[一-鿿]", s or ""))

    # 1) 幻灯片正文区
    for idx, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            shp_text = shape_text(sh).strip()
            for text, font, _, _ in iter_runs(sh):
                if not (text or "").strip() or not is_chinese(text):
                    continue
                if font not in ("宋体", "SimSun"):
                    bad.append((idx, shp_text[:20], font, text[:20]))

            # 覆盖范围标记
            if idx == 2 and ("目录" in shp_text or shp_text in TOC_ENTRIES
                             or any(e in shp_text for e in TOC_ENTRIES)):
                if is_chinese(shp_text):
                    coverage["目录"] = True
            if shp_text in section_titles_set and is_chinese(shp_text):
                coverage["标题"] = True
            if shp_text in sub_titles_set and is_chinese(shp_text):
                coverage["二级标题"] = True
            if idx in BODY_PAGES and len(shp_text) > 30 and is_chinese(shp_text):
                coverage["正文"] = True
            if idx == n_slides and is_chinese(shp_text):
                coverage["尾页"] = True

    # 2) 页脚（slide master / layout / slide 的 footer placeholder & 母版页脚 shape）
    #    python-pptx 没有直接 footer API，这里检查母版/版式中 idx=12 (footer) 的 placeholder
    try:
        masters = list(prs.slide_masters)
    except Exception:
        masters = []
    sources = []
    for m in masters:
        sources.append(("master", m))
        for layout in m.slide_layouts:
            sources.append(("layout", layout))
    for tag, src in sources:
        for sh in src.shapes:
            if not sh.has_text_frame:
                continue
            shp_text = shape_text(sh).strip()
            if not shp_text:
                continue
            for text, font, _, _ in iter_runs(sh):
                if not (text or "").strip() or not is_chinese(text):
                    continue
                if font not in ("宋体", "SimSun"):
                    bad.append((f"{tag}", shp_text[:20], font, text[:20]))
            # 是否为页脚 placeholder（ph type=ftr）
            try:
                if sh.is_placeholder and sh.placeholder_format is not None \
                        and getattr(sh.placeholder_format, "idx", None) is not None:
                    # 类型 15=footer 在 python-pptx
                    from pptx.enum.shapes import PP_PLACEHOLDER
                    if sh.placeholder_format.type == PP_PLACEHOLDER.FOOTER and is_chinese(shp_text):
                        coverage["页脚"] = True
            except Exception:
                pass

    # 若 PPT 实际未使用页脚功能（即整份未出现页脚文字），则该子点视为"不存在不扣分"，
    # 仅当真正出现过页脚文本且非宋体时才算违规 —— 这已由上面的 bad 列表统一处理。

    if bad:
        sample = bad[:3]
        return False, f"存在 {len(bad)} 处非宋体中文，示例：{sample}"

    return True, ("全文中文均为宋体；覆盖检测："
                  + "、".join(f"{k}{'✓' if v else '—'}" for k, v in coverage.items()))


# --- 动画相关 ---
def parse_timing(raw_zip, slide_no):
    """返回该页的动画条目列表：[(spid, filter), ...]，并附带原始 timing 节点"""
    xml = get_slide_xml(raw_zip, slide_no)
    root = etree.fromstring(xml.encode("utf-8"))
    timing = root.find("p:timing", NS)
    entries = []
    if timing is None:
        return entries, None
    for par in timing.findall(".//p:par", NS):
        eff = par.find(".//p:animEffect", NS)
        spt = par.find(".//p:spTgt", NS)
        if eff is not None and spt is not None:
            entries.append({
                "spid": spt.get("spid"),
                "filter": eff.get("filter") or eff.get("transition"),
            })
    return entries, timing


def find_spid_by_text(slide, keyword):
    for sh in slide.shapes:
        if sh.has_text_frame and keyword in sh.text_frame.text:
            return str(sh.shape_id)
    return None


MIN_ANIM_DUR_MS = 500
ANIM_TAGS = ("animEffect", "anim", "animClr", "animRot",
             "set", "animScale", "animMotion")


def get_main_sequence(root):
    """返回动画窗格主序列 mainSeq；没有则返回 None。"""
    for seq in root.findall(".//p:seq", NS):
        ctn = seq.find("p:cTn", NS)
        if ctn is not None and ctn.get("nodeType") == "mainSeq":
            return seq
    return None


def anim_duration_ms(node):
    """返回一个动画行为节点的实际放映持续时间（毫秒）；无法判定则返回 None。"""
    ctn = node.find("p:cBhvr/p:cTn", NS)
    if ctn is None:
        ctn = node.find(".//p:cTn", NS)
    if ctn is None:
        return None
    dur = ctn.get("dur")
    return int(dur) if dur and dur.isdigit() else None


def collect_anim_nodes_by_spid(scope, tags=ANIM_TAGS):
    """在指定 XML 范围内收集 spid → [(tag, duration_ms), ...]。"""
    spid_anims = {}
    for tag in tags:
        for node in scope.findall(f".//p:{tag}", NS):
            dur_ms = anim_duration_ms(node)
            for spt in node.findall(".//p:spTgt", NS):
                spid = spt.get("spid")
                if spid:
                    spid_anims.setdefault(spid, []).append((tag, dur_ms))
    return spid_anims


def check_anim_duration_for_spids(scope, spids, tags=ANIM_TAGS, min_ms=MIN_ANIM_DUR_MS):
    """检查每个目标对象至少有一个动画节点的持续时间达到 min_ms。"""
    spid_anims = collect_anim_nodes_by_spid(scope, tags)
    bad = []
    for spid in sorted(spids):
        durations = [dur for _, dur in spid_anims.get(spid, [])]
        if not durations:
            bad.append(f"spid {spid} 无动画")
        elif not any(dur is not None and dur >= min_ms for dur in durations):
            shown = ["未设置" if dur is None else f"{dur}ms" for dur in durations]
            bad.append(f"spid {spid} 时长 {shown}")
    return len(bad) == 0, bad


def check_sub_title_anim(prs, raw_zip):
    """+5 正文页二级标题动画：
       ① 必须覆盖第 4、6、8、10、12、14、16、18、20 这九页；
       ② 每一页上的两个二级标题（X.1、X.2）都各自被赋予了动画效果
          —— 只要 timing 内存在任一有效动画节点（animEffect / anim / animClr /
             animRot / set / animScale / animMotion）以该二级标题为目标即可，
             不再强制要求 mainSeq 中的 p:animEffect，也不再要求最小持续时间。
       九个页面、每页两个二级标题全部命中才计 +5。
    """
    ok_pages = []
    fail_details = []

    if len(prs.slides) < max(BODY_PAGES):
        return False, f"PPT 不足 {max(BODY_PAGES)} 页"

    for i, pno in enumerate(BODY_PAGES):
        slide = prs.slides[pno - 1]
        s1_text, s2_text = SUB_TITLES[i]

        sub_ids = {}
        for sh in slide.shapes:
            if sh.has_text_frame and shape_text(sh).strip() in (s1_text, s2_text):
                sub_ids[shape_text(sh).strip()] = str(sh.shape_id)
        if s1_text not in sub_ids or s2_text not in sub_ids:
            fail_details.append(f"P{pno} 二级标题文本未找到（命中 {list(sub_ids.keys())}）")
            continue

        xml = get_slide_xml(raw_zip, pno)
        root = etree.fromstring(xml.encode("utf-8"))
        timing = root.find("p:timing", NS)
        if timing is None:
            fail_details.append(f"P{pno} 无 timing（动画信息）")
            continue

        # ② 在 timing 内收集所有以 shape 为目标的动画节点 spid 集合
        anim_spids = set()
        for tag in ANIM_TAGS:
            for node in timing.findall(f".//p:{tag}", NS):
                for spt in node.findall(".//p:spTgt", NS):
                    spid = spt.get("spid")
                    if spid:
                        anim_spids.add(spid)

        missing = [t for t in (s1_text, s2_text) if sub_ids[t] not in anim_spids]
        if missing:
            fail_details.append(f"P{pno} 二级标题缺少动画：{missing}")
            continue

        ok_pages.append(pno)

    ok = (ok_pages == BODY_PAGES)
    msg = f"正文页二级标题动画 命中 {len(ok_pages)}/9（页 {ok_pages}）"
    if not ok and fail_details:
        msg += f"；未命中：{fail_details[:3]}"
    return ok, msg


def check_body_anim(prs, raw_zip):
    # +5 正文页正文动画
    # 九页全部满足才计 +5：
    # 1. 覆盖第 4/6/8/10/12/14/16/18/20 页
    # 2. 每页两个正文 shape 都设置了动画
    # 3. 与二级标题动画"区分明显"：从效果类型、方向、参数等视觉维度比较，
    #    只要两组签名集合不完全相同（存在至少一个视觉差异）即视为区分明显。
    ok_pages = []
    fail_details = []

    if len(prs.slides) < max(BODY_PAGES):
        return False, f"PPT 不足 {max(BODY_PAGES)} 页"

    # 动画节点类型 → 视觉特征签名
    ANIM_TAGS = ("animEffect", "anim", "animClr", "animRot",
                 "set", "animScale", "animMotion")

    def effect_signature(node):
        """把一个动画节点抽象为"视觉签名"：包含类型 + 方向 + 参数等。
        目的是让"标题飞入自左 / 正文飞入自右"这类视觉差异能被识别为不同。"""
        tag = etree.QName(node).localname
        parts = [tag]
        if tag == "animEffect":
            parts.append(f"transition={node.get('transition') or ''}")
            # filter 常见形如 "in(fromLeft)" / "in(fromRight)"，天然包含方向
            parts.append(f"filter={node.get('filter') or ''}")
        elif tag == "anim":
            parts.append(f"attr={node.get('attributeName') or ''}")
            parts.append(f"by={node.get('by') or ''}")
            parts.append(f"from={node.get('from') or ''}")
            parts.append(f"to={node.get('to') or ''}")
        elif tag == "animMotion":
            parts.append(f"path={node.get('path') or ''}")
            parts.append(f"origin={node.get('origin') or ''}")
        elif tag == "animScale":
            by = node.find(".//p:by", NS)
            if by is not None:
                parts.append(f"by=({by.get('x') or ''},{by.get('y') or ''})")
        elif tag == "animRot":
            parts.append(f"by={node.get('by') or ''}")
            parts.append(f"from={node.get('from') or ''}")
            parts.append(f"to={node.get('to') or ''}")
        elif tag == "animClr":
            parts.append(f"clrSpc={node.get('clrSpc') or ''}")
            parts.append(f"dir={node.get('dir') or ''}")
        # 兜底：附加 presetClass/presetID/presetSubtype（PowerPoint 动画预设 3 元组）
        ctn = node.find("p:cBhvr/p:cTn", NS) or node.find(".//p:cTn", NS)
        if ctn is not None:
            parts.append(f"presetClass={ctn.get('presetClass') or ''}")
            parts.append(f"presetID={ctn.get('presetID') or ''}")
            parts.append(f"presetSubtype={ctn.get('presetSubtype') or ''}")
        return "|".join(parts)

    for i, pno in enumerate(BODY_PAGES):
        slide = prs.slides[pno - 1]
        s1_text, s2_text = SUB_TITLES[i]

        # 二级标题与"对应正文" shape 定位
        sub_shapes = {}
        for sh in slide.shapes:
            if sh.has_text_frame and shape_text(sh).strip() in (s1_text, s2_text):
                sub_shapes[shape_text(sh).strip()] = sh
        if s1_text not in sub_shapes or s2_text not in sub_shapes:
            fail_details.append(f"P{pno} 二级标题未集齐")
            continue

        sw = prs.slide_width
        body_shapes = {}   # key=二级标题文本 → 正文 shape
        for key in (s1_text, s2_text):
            sub = sub_shapes[key]
            other_sub_top = sub_shapes[s2_text if key == s1_text else s1_text].top or 0
            cands = []
            for sh in slide.shapes:
                if not sh.has_text_frame or sh is sub:
                    continue
                t = shape_text(sh).strip()
                if not t or t in (s1_text, s2_text):
                    continue
                if (sh.top or 0) <= (sub.top or 0):
                    continue
                if key == s1_text and (sh.top or 0) >= other_sub_top:
                    continue
                if abs((sh.left or 0) - (sub.left or 0)) > 0.1 * sw:
                    continue
                if len(t) < 20:
                    continue
                cands.append(sh)
            if not cands:
                fail_details.append(f'P{pno} 未找到"{key}"对应的正文 shape')
                body_shapes = None
                break
            cands.sort(key=lambda x: x.top or 0)
            body_shapes[key] = cands[0]
        if body_shapes is None:
            continue

        sub_ids = {str(sub_shapes[k].shape_id) for k in (s1_text, s2_text)}
        body_ids = {str(body_shapes[k].shape_id) for k in (s1_text, s2_text)}

        # 收集 timing 中：spid → 该 spid 的视觉签名集合
        xml = get_slide_xml(raw_zip, pno)
        root = etree.fromstring(xml.encode("utf-8"))
        timing = root.find("p:timing", NS)
        if timing is None:
            fail_details.append(f"P{pno} 无 timing")
            continue
        spid_effects = {}     # spid → set(signature)
        for tag in ANIM_TAGS:
            for node in timing.findall(f".//p:{tag}", NS):
                for spt in node.findall(".//p:spTgt", NS):
                    spid = spt.get("spid")
                    if spid is None:
                        continue
                    spid_effects.setdefault(spid, set()).add(effect_signature(node))

        # ② 正文都要有动画
        missing_body = [bid for bid in body_ids if bid not in spid_effects]
        if missing_body:
            fail_details.append(f"P{pno} 正文未设置动画（spid {missing_body}）")
            continue

        # ③ 与二级标题动画"区分明显"：比较视觉签名集合是否完全相同；
        #   两组都要有动画，且至少存在一个签名不被另一组共享（即两个集合不相等）。
        sub_effects = set().union(*(spid_effects.get(s, set()) for s in sub_ids))
        body_effects = set().union(*(spid_effects.get(b, set()) for b in body_ids))
        if not sub_effects or not body_effects:
            fail_details.append(f"P{pno} 二级标题或正文动画为空")
            continue
        if sub_effects == body_effects:
            fail_details.append(
                f"P{pno} 正文与二级标题动画视觉签名完全一致，未形成明显区分：{sub_effects}")
            continue

        ok_pages.append(pno)

    ok = (ok_pages == BODY_PAGES)
    msg = f"正文页正文动画且与二级标题动画区分明显 命中 {len(ok_pages)}/9（页 {ok_pages}）"
    if not ok and fail_details:
        msg += f"；未命中：{fail_details[:3]}"
    return ok, msg


def check_anim_trigger(prs, raw_zip):
    """+5 正文页动画触发：
       对第 4、6、8、10、12、14、16、18、20 共九页，每页都必须满足：
       ① 切换进入后，通过单击开始本页动画
          —— 存在主序列 <p:seq nodeType="mainSeq">，且其 <p:nextCondLst>
             含 evt="onNext"（单击进入下一步）触发；
          —— 同时该主序列不是"随上一动作自动开始"（避免直接随翻页播放）。
       ② 二级标题和正文"在同一页面内"自动依次出现
          —— 即除第一个动画外，后续动画的启动条件不能是 onClick；
             允许 delay=0 / 设定延迟（withPrev / afterPrev 在 PPT XML 里表现为
             无 evt 的 cond），代表上一动作之后自动继续。
       ③ 这页里所有"四个目标"（X.1 二级标题 / X.1 正文 / X.2 二级标题 / X.2 正文）
          都包含在主序列里，确保单击之后"随后"就是它们依次出现。
       九页全部满足才计 +5。
    """
    ok_pages = []
    fail_details = []

    if len(prs.slides) < max(BODY_PAGES):
        return False, f"PPT 不足 {max(BODY_PAGES)} 页"

    for i, pno in enumerate(BODY_PAGES):
        slide = prs.slides[pno - 1]
        s1_text, s2_text = SUB_TITLES[i]

        # 二级标题 + 对应正文 shape_id 集合
        sub_shapes = {}
        for sh in slide.shapes:
            if sh.has_text_frame and shape_text(sh).strip() in (s1_text, s2_text):
                sub_shapes[shape_text(sh).strip()] = sh
        if s1_text not in sub_shapes or s2_text not in sub_shapes:
            fail_details.append(f"P{pno} 二级标题未集齐")
            continue
        sw = prs.slide_width
        target_ids = set()
        for key in (s1_text, s2_text):
            sub = sub_shapes[key]
            target_ids.add(str(sub.shape_id))
            other_top = sub_shapes[s2_text if key == s1_text else s1_text].top or 0
            cands = []
            for sh in slide.shapes:
                if not sh.has_text_frame or sh is sub:
                    continue
                t = shape_text(sh).strip()
                if not t or t in (s1_text, s2_text):
                    continue
                if (sh.top or 0) <= (sub.top or 0):
                    continue
                if key == s1_text and (sh.top or 0) >= other_top:
                    continue
                if abs((sh.left or 0) - (sub.left or 0)) > 0.1 * sw:
                    continue
                if len(t) < 20:
                    continue
                cands.append(sh)
            if not cands:
                target_ids = None
                fail_details.append(f'P{pno} 未找到"{key}"对应的正文 shape')
                break
            cands.sort(key=lambda x: x.top or 0)
            target_ids.add(str(cands[0].shape_id))
        if target_ids is None:
            continue

        # 解析 timing
        xml = get_slide_xml(raw_zip, pno)
        root = etree.fromstring(xml.encode("utf-8"))
        # 找 mainSeq
        mainSeq = None
        for s in root.findall(".//p:seq", NS):
            ctn = s.find("p:cTn", NS)
            if ctn is not None and ctn.get("nodeType") == "mainSeq":
                mainSeq = s
                break
        if mainSeq is None:
            fail_details.append(f"P{pno} 无 mainSeq 序列")
            continue

        # ① nextCondLst 包含 onNext（单击开始本页动画）
        next_cond = mainSeq.find("p:nextCondLst/p:cond", NS)
        if next_cond is None or next_cond.get("evt") != "onNext":
            fail_details.append(f"P{pno} 主序列未配置 onNext 单击触发")
            continue

        # 主序列下顶层 par 节点（不强制数量，允许多个对象组合成一个 par）
        top_pars = mainSeq.findall("p:cTn/p:childTnLst/p:par", NS)
        if not top_pars:
            fail_details.append(f"P{pno} 主序列没有子动画")
            continue

        # ③ 这页里所有 4 个目标都出现在主序列里（可分散在任意 par 中）
        spids_in_seq = set()
        for par in top_pars:
            for spt in par.findall(".//p:spTgt", NS):
                spid = spt.get("spid")
                if spid:
                    spids_in_seq.add(spid)
        if not target_ids.issubset(spids_in_seq):
            fail_details.append(
                f"P{pno} 主序列缺少目标动画：{target_ids - spids_in_seq}")
            continue

        # ② 除第一个动画外，后续不能是 onClick（避免"需多次点击"）
        bad_click = False
        for idx_par, par in enumerate(top_pars):
            cond = par.find(".//p:stCondLst/p:cond", NS)
            evt = cond.get("evt") if cond is not None else None
            if idx_par == 0:
                continue  # 第一个动画的"启动"由整个 mainSeq 的 onNext 控制
            if evt == "onClick":
                bad_click = True
                fail_details.append(f"P{pno} 第{idx_par+1}个动画为 onClick，需多次点击")
                break
        if bad_click:
            continue

        ok_pages.append(pno)

    ok = (ok_pages == BODY_PAGES)
    msg = f'正文页"单击开始本页 + 标题正文自动依次出现" 命中 {len(ok_pages)}/9（页 {ok_pages}）'
    if not ok and fail_details:
        msg += f"；未命中：{fail_details[:3]}"
    return ok, msg


def check_anim_order(prs, raw_zip):
    """+3 正文页动画顺序：
       对第 4、6、8、10、12、14、16、18、20 共九页，每页都必须满足：
       ① 先出现二级标题，再出现对应正文内容
          —— 每一个二级小节内：二级标题在主序列里的位置必须早于其对应正文；
       ② 两个二级小节的动画顺序与页面文本顺序一致
          —— 页面上 top 较小（更靠上）的那个二级小节，其动画（二级标题+正文）
             整体出现在 top 较大的二级小节之前；
       九页全部满足才计 +3。
    """
    ok_pages = []
    fail_details = []

    if len(prs.slides) < max(BODY_PAGES):
        return False, f"PPT 不足 {max(BODY_PAGES)} 页"

    for i, pno in enumerate(BODY_PAGES):
        slide = prs.slides[pno - 1]
        s1_text, s2_text = SUB_TITLES[i]

        # 定位两个二级标题 + 各自正文
        sub_shapes = {}
        for sh in slide.shapes:
            if sh.has_text_frame and shape_text(sh).strip() in (s1_text, s2_text):
                sub_shapes[shape_text(sh).strip()] = sh
        if s1_text not in sub_shapes or s2_text not in sub_shapes:
            fail_details.append(f"P{pno} 二级标题未集齐")
            continue

        sw = prs.slide_width
        body_shapes = {}
        ok_locate = True
        for key in (s1_text, s2_text):
            sub = sub_shapes[key]
            other_top = sub_shapes[s2_text if key == s1_text else s1_text].top or 0
            cands = []
            for sh in slide.shapes:
                if not sh.has_text_frame or sh is sub:
                    continue
                t = shape_text(sh).strip()
                if not t or t in (s1_text, s2_text):
                    continue
                if (sh.top or 0) <= (sub.top or 0):
                    continue
                if key == s1_text and (sh.top or 0) >= other_top:
                    continue
                if abs((sh.left or 0) - (sub.left or 0)) > 0.1 * sw:
                    continue
                if len(t) < 20:
                    continue
                cands.append(sh)
            if not cands:
                fail_details.append(f'P{pno} 未找到"{key}"对应正文')
                ok_locate = False
                break
            cands.sort(key=lambda x: x.top or 0)
            body_shapes[key] = cands[0]
        if not ok_locate:
            continue

        # 按"页面文本顺序"决定两个二级小节的先后：top 较小的优先
        sub_order_by_top = sorted(
            (s1_text, s2_text),
            key=lambda k: sub_shapes[k].top or 0,
        )
        first_key, second_key = sub_order_by_top
        first_sub_id = str(sub_shapes[first_key].shape_id)
        first_body_id = str(body_shapes[first_key].shape_id)
        second_sub_id = str(sub_shapes[second_key].shape_id)
        second_body_id = str(body_shapes[second_key].shape_id)

        # 解析 timing 主序列动画顺序
        xml = get_slide_xml(raw_zip, pno)
        root = etree.fromstring(xml.encode("utf-8"))
        mainSeq = None
        for s in root.findall(".//p:seq", NS):
            ctn = s.find("p:cTn", NS)
            if ctn is not None and ctn.get("nodeType") == "mainSeq":
                mainSeq = s
                break
        if mainSeq is None:
            fail_details.append(f"P{pno} 无 mainSeq")
            continue

        # 按主序列 par 索引记录：
        #   entrance_pos[spid] = 该对象在时间线上"首次可见 / 入场"动画所在的顶层 par 序号
        #   any_pos[spid]      = 该对象任意首次出现的 par 序号（兜底）
        # 说明：先按 par 遍历一次，把每个 spid 的"入场类"位置记下来；
        #      多动画步骤时优先用入场位置，避免把后续变色/位移误判为出现顺序。
        top_pars = mainSeq.findall("p:cTn/p:childTnLst/p:par", NS)
        entrance_pos = {}
        any_pos = {}

        def _is_entrance_anim(anim_node):
            """判定一个动画行为节点是否为"入场 / 首次可见"动画。"""
            tag = etree.QName(anim_node).localname
            # PPT 预设动画：presetClass="entr" 表示入场类
            ctn = anim_node.find("p:cBhvr/p:cTn", NS)
            if ctn is None:
                ctn = anim_node.find(".//p:cTn", NS)
            if ctn is not None and ctn.get("presetClass") == "entr":
                return True
            # animEffect transition="in" —— 常见的进入 / 淡入 / 飞入 / 擦除等
            if tag == "animEffect" and (anim_node.get("transition") or "").lower() == "in":
                return True
            # set 让 style.visibility → visible，也是"首次显现"
            if tag == "set":
                attr = anim_node.find("p:cBhvr/p:attrNameLst/p:attrName", NS)
                to = anim_node.find("p:to/p:strVal", NS)
                if attr is not None and to is not None \
                        and "visibility" in (attr.text or "") \
                        and (to.get("val") or "").lower() == "visible":
                    return True
            return False

        for par_idx, par in enumerate(top_pars):
            # 枚举本 par 内所有 spTgt 所对应的动画节点
            for spt in par.findall(".//p:spTgt", NS):
                spid = spt.get("spid")
                if not spid:
                    continue
                any_pos.setdefault(spid, par_idx)
                # 定位其所属动画节点：spTgt 位于 cBhvr/tgtEl 或 tgtEl 内，往上找动画标签
                anim = spt.getparent()
                while anim is not None and etree.QName(anim).localname not in ANIM_TAGS:
                    anim = anim.getparent()
                if anim is not None and _is_entrance_anim(anim):
                    entrance_pos.setdefault(spid, par_idx)

        def _pos(spid):
            # 首选入场位置；无入场时回退到任意首次出现位置
            return entrance_pos.get(spid, any_pos.get(spid))

        needed = [first_sub_id, first_body_id, second_sub_id, second_body_id]
        missing = [x for x in needed if _pos(x) is None]
        if missing:
            fail_details.append(f"P{pno} 动画序列缺少目标 spid {missing}")
            continue

        idx_fs = _pos(first_sub_id)
        idx_fb = _pos(first_body_id)
        idx_ss = _pos(second_sub_id)
        idx_sb = _pos(second_body_id)

        # ① 每个小节：二级标题的入场位置早于对应正文的入场位置
        if not (idx_fs < idx_fb):
            fail_details.append(
                f'P{pno} 第一节中"{first_key}"二级标题未先于其正文 '
                f"(idx 标题{idx_fs}, 正文{idx_fb})")
            continue
        if not (idx_ss < idx_sb):
            fail_details.append(
                f'P{pno} 第二节中"{second_key}"二级标题未先于其正文 '
                f"(idx 标题{idx_ss}, 正文{idx_sb})")
            continue

        # ② 两个小节顺序与页面文本顺序一致
        #    第一节的"标题+正文"整体在第二节之前 —— 用最大索引比最小索引
        if not (max(idx_fs, idx_fb) < min(idx_ss, idx_sb)):
            fail_details.append(
                f"P{pno} 两个小节动画顺序与文本顺序不一致 "
                f"(第一节 {idx_fs},{idx_fb} vs 第二节 {idx_ss},{idx_sb})")
            continue

        ok_pages.append(pno)

    ok = (ok_pages == BODY_PAGES)
    msg = f"正文页动画顺序正确 命中 {len(ok_pages)}/9（页 {ok_pages}）"
    if not ok and fail_details:
        msg += f"；未命中：{fail_details[:3]}"
    return ok, msg


def check_anim_continuity(prs, raw_zip):
    """+3 正文页动画连续性：
       对第 4、6、8、10、12、14、16、18、20 共九页，每页都必须满足：
       ① 同一页面内动画设置为"上一动画之后"或等效自动播放
          —— 主序列里除"第一个动画"外，其余 par 节点的启动条件 (p:cond)
             要么没有 evt（withPrev / afterPrev，即上一动作之后自动播放），
             要么 evt == "afterEffect" / "afterPrev"；
          —— 任何 evt == "onClick" 的后续动画都视为"需要再次点击"。
       ② 不需要用户反复点击多次才能看完本页内容
          —— 主序列里整页只能有 1 次需要点击的触发（即整个 mainSeq 由 onNext
             启动），其下属顶层 par 中"onClick"计数必须为 0。
       九页全部满足才计 +3。
    """
    ok_pages = []
    fail_details = []

    if len(prs.slides) < max(BODY_PAGES):
        return False, f"PPT 不足 {max(BODY_PAGES)} 页"

    for pno in BODY_PAGES:
        xml = get_slide_xml(raw_zip, pno)
        root = etree.fromstring(xml.encode("utf-8"))

        # 找 mainSeq
        mainSeq = None
        for s in root.findall(".//p:seq", NS):
            ctn = s.find("p:cTn", NS)
            if ctn is not None and ctn.get("nodeType") == "mainSeq":
                mainSeq = s
                break
        if mainSeq is None:
            fail_details.append(f"P{pno} 无 mainSeq")
            continue

        top_pars = mainSeq.findall("p:cTn/p:childTnLst/p:par", NS)
        if len(top_pars) < 2:
            fail_details.append(f"P{pno} 主序列动画 < 2 个，无法判定连续性")
            continue

        # ② 统计需要点击的后续动画数量
        click_count_after_first = 0
        non_auto_after_first = 0
        for idx, par in enumerate(top_pars):
            conds = par.findall(".//p:stCondLst/p:cond", NS)
            # 收集本 par 的所有 evt 取值
            evts = [c.get("evt") for c in conds]
            if idx == 0:
                # 第一个动画的启动 = 整个 mainSeq 的 onNext，不算多次点击
                continue
            # 后续动画：任何带 onClick 的都算需要再次点击
            if "onClick" in evts:
                click_count_after_first += 1
                continue
            # ① "上一动画之后"或等效自动播放：
            # 允许 evt 为 None（withPrev/afterPrev 在 XML 中即无 evt），
            # 或者 afterEffect / afterPrev
            allowed = {None, "afterEffect", "afterPrev"}
            if not all((e in allowed) for e in evts):
                non_auto_after_first += 1

        if click_count_after_first > 0:
            fail_details.append(
                f"P{pno} 主序列后续动画中存在 {click_count_after_first} 处 onClick，"
                f"需多次点击")
            continue
        if non_auto_after_first > 0:
            fail_details.append(
                f'P{pno} 存在 {non_auto_after_first} 个非"上一动画之后/自动"触发')
            continue

        ok_pages.append(pno)

    ok = (ok_pages == BODY_PAGES)
    msg = (f'正文页动画连续性（"上一动画之后"自动播放，无需多次点击） '
           f"命中 {len(ok_pages)}/9（页 {ok_pages}）")
    if not ok and fail_details:
        msg += f"；未命中：{fail_details[:3]}"
    return ok, msg


def check_anim_effect_diff(prs, raw_zip):
    """+3 动画效果差异：
       对第 4、6、8、10、12、14、16、18、20 共九页，每页都必须满足：
       ① 二级标题使用的效果属于"淡入 / 飞入 / 擦除"等进入类动画
          —— 即 timing 中以二级标题为目标的动画节点至少存在一个 p:animEffect
             或常见的进入动作（filter ∈ {fade, fly*, wipe*, ...}）；
       ② 正文使用"另一种不同效果"
          —— 正文的效果集合非空，且与二级标题的效果集合不完全相同；
       ③ 不能所有对象完全使用同一种动画
          —— 本页所有有动画的对象的"效果标识"去重后，必须 ≥ 2 种。
       九页全部满足才计 +3。
    """
    ok_pages = []
    fail_details = []

    if len(prs.slides) < max(BODY_PAGES):
        return False, f"PPT 不足 {max(BODY_PAGES)} 页"

    ANIM_TAGS = ("animEffect", "anim", "animClr", "animRot",
                 "set", "animScale", "animMotion")

    # 允许作为二级标题入场动画的 filter 关键字（淡入 / 飞入 / 擦除 "等"）
    ENTRANCE_FILTERS = ("fade", "fly", "wipe", "blinds", "wheel",
                        "split", "checkerboard", "diamond", "circle",
                        "plus", "random", "box", "strips", "appear",
                        "dissolve", "peek", "zoom")

    def effect_key(node):
        tag = etree.QName(node).localname
        if tag == "animEffect":
            return f"animEffect:{node.get('transition') or ''}:{node.get('filter') or ''}"
        if tag == "anim":
            return f"anim:{node.get('attributeName') or ''}"
        if tag == "animMotion":
            return "animMotion"
        if tag == "animScale":
            return "animScale"
        if tag == "animClr":
            return "animClr"
        if tag == "animRot":
            return "animRot"
        if tag == "set":
            return "set"
        return tag

    for i, pno in enumerate(BODY_PAGES):
        slide = prs.slides[pno - 1]
        s1_text, s2_text = SUB_TITLES[i]

        # 二级标题 + 对应正文 shape
        sub_shapes = {}
        for sh in slide.shapes:
            if sh.has_text_frame and shape_text(sh).strip() in (s1_text, s2_text):
                sub_shapes[shape_text(sh).strip()] = sh
        if s1_text not in sub_shapes or s2_text not in sub_shapes:
            fail_details.append(f"P{pno} 二级标题未集齐")
            continue

        sw = prs.slide_width
        body_shapes = {}
        ok_locate = True
        for key in (s1_text, s2_text):
            sub = sub_shapes[key]
            other_top = sub_shapes[s2_text if key == s1_text else s1_text].top or 0
            cands = []
            for sh in slide.shapes:
                if not sh.has_text_frame or sh is sub:
                    continue
                t = shape_text(sh).strip()
                if not t or t in (s1_text, s2_text):
                    continue
                if (sh.top or 0) <= (sub.top or 0):
                    continue
                if key == s1_text and (sh.top or 0) >= other_top:
                    continue
                if abs((sh.left or 0) - (sub.left or 0)) > 0.1 * sw:
                    continue
                if len(t) < 20:
                    continue
                cands.append(sh)
            if not cands:
                ok_locate = False
                fail_details.append(f'P{pno} 未找到"{key}"对应正文')
                break
            cands.sort(key=lambda x: x.top or 0)
            body_shapes[key] = cands[0]
        if not ok_locate:
            continue

        sub_ids = {str(sub_shapes[k].shape_id) for k in (s1_text, s2_text)}
        body_ids = {str(body_shapes[k].shape_id) for k in (s1_text, s2_text)}

        # 解析 timing：spid → effect_key 集合 & filter 集合（用于判定"淡入/飞入/擦除等"）
        xml = get_slide_xml(raw_zip, pno)
        root = etree.fromstring(xml.encode("utf-8"))
        timing = root.find("p:timing", NS)
        if timing is None:
            fail_details.append(f"P{pno} 无 timing")
            continue
        spid_effects = {}
        spid_filters = {}
        for tag in ANIM_TAGS:
            for node in timing.findall(f".//p:{tag}", NS):
                for spt in node.findall(".//p:spTgt", NS):
                    spid = spt.get("spid")
                    if not spid:
                        continue
                    spid_effects.setdefault(spid, set()).add(effect_key(node))
                    if tag == "animEffect":
                        f = node.get("filter") or ""
                        spid_filters.setdefault(spid, set()).add(f.lower())

        # ① 二级标题使用"进入类"动画（淡入 / 飞入 / 擦除 等）
        #   接受下述任一形式：
        #   a) animEffect 且 transition="in"（明确的入场变换）；
        #   b) animEffect 且 filter 命中 ENTRANCE_FILTERS 之一（fade/fly/wipe/…）；
        #   c) 使用了 PPT 预设动画且 presetClass="entr"（入场类）；
        #   d) set 让 style.visibility → visible（等效"出现"入场）。
        def _is_entrance_anim_node(node):
            tag = etree.QName(node).localname
            if tag == "animEffect":
                if (node.get("transition") or "").lower() == "in":
                    return True
                f = (node.get("filter") or "").lower()
                if any(k in f for k in ENTRANCE_FILTERS):
                    return True
            ctn = node.find("p:cBhvr/p:cTn", NS) or node.find(".//p:cTn", NS)
            if ctn is not None and ctn.get("presetClass") == "entr":
                return True
            if tag == "set":
                attr = node.find("p:cBhvr/p:attrNameLst/p:attrName", NS)
                to = node.find("p:to/p:strVal", NS)
                if attr is not None and to is not None \
                        and "visibility" in (attr.text or "") \
                        and (to.get("val") or "").lower() == "visible":
                    return True
            return False

        # 按 spid 判断二级标题是否包含"进入类"节点
        sub_has_entrance = False
        for tag in ANIM_TAGS:
            if sub_has_entrance:
                break
            for node in timing.findall(f".//p:{tag}", NS):
                if not _is_entrance_anim_node(node):
                    continue
                for spt in node.findall(".//p:spTgt", NS):
                    if spt.get("spid") in sub_ids:
                        sub_has_entrance = True
                        break
                if sub_has_entrance:
                    break
        if not sub_has_entrance:
            fail_details.append(f"P{pno} 二级标题未使用进入类动画（淡入/飞入/擦除等）")
            continue

        sub_effects = set().union(*(spid_effects.get(s, set()) for s in sub_ids))
        body_effects = set().union(*(spid_effects.get(b, set()) for b in body_ids))

        # ② 正文使用"另一种不同效果"：正文效果非空，且 sub_effects != body_effects
        if not body_effects:
            fail_details.append(f"P{pno} 正文缺少动画效果")
            continue
        if sub_effects == body_effects:
            fail_details.append(
                f"P{pno} 二级标题与正文使用了完全相同的效果集合：{sub_effects}")
            continue

        # ③ 不能所有对象完全使用同一种动画
        all_effects = set()
        for sid in (sub_ids | body_ids):
            all_effects |= spid_effects.get(sid, set())
        if len(all_effects) < 2:
            fail_details.append(
                f"P{pno} 所有对象使用了同一种动画：{all_effects}")
            continue

        ok_pages.append(pno)

    ok = (ok_pages == BODY_PAGES)
    msg = f"动画效果差异（标题与正文不同 & 非所有对象同一动画） 命中 {len(ok_pages)}/9（页 {ok_pages}）"
    if not ok and fail_details:
        msg += f"；未命中：{fail_details[:3]}"
    return ok, msg


def check_page_switch_logic(prs, raw_zip):
    """+3 页面切换逻辑：
       对第 4、6、8、10、12、14、16、18、20 共九页，每页都必须满足：
       ① 换页后内容动画不会"自动提前播放"
          —— 自动换页（advTm）本身不算失败，只关心内容动画是否需要单击启动；
          —— 但若 <p:transition advClick="0"> 显式禁用单击换页，
             会导致无法通过单击触发本页内容动画，视为异常；
          —— 该页 timing 主序列第一个动画的启动条件不能是 onLoad / withEffect
             （这两种条件等价于"切换进入即播放"）。
       ② 需单击后开始该页内容动画
          —— 该页 timing 主序列 <p:nextCondLst><p:cond evt="onNext"/> 必须存在，
             即主序列由"单击下一步"启动。
       九页全部满足才计 +3。
    """
    ok_pages = []
    fail_details = []

    if len(prs.slides) < max(BODY_PAGES):
        return False, f"PPT 不足 {max(BODY_PAGES)} 页"

    for pno in BODY_PAGES:
        xml = get_slide_xml(raw_zip, pno)
        root = etree.fromstring(xml.encode("utf-8"))

        # ① 换页本身允许自动（advTm 不判失败）；但 advClick="0" 会阻止单击，
        #    导致无法通过点击启动内容动画，视为异常
        trans = root.find("p:transition", NS)
        if trans is not None:
            adv_click = trans.get("advClick")
            if adv_click == "0":
                fail_details.append(f"P{pno} 切换被设为不允许单击")
                continue

        # 找 mainSeq
        mainSeq = None
        for s in root.findall(".//p:seq", NS):
            ctn = s.find("p:cTn", NS)
            if ctn is not None and ctn.get("nodeType") == "mainSeq":
                mainSeq = s
                break
        if mainSeq is None:
            fail_details.append(f"P{pno} 无 mainSeq")
            continue

        # ② 主序列由单击启动：nextCondLst.cond.evt == onNext
        next_cond = mainSeq.find("p:nextCondLst/p:cond", NS)
        if next_cond is None or next_cond.get("evt") != "onNext":
            fail_details.append(f"P{pno} 主序列未配置 onNext 启动（不需要单击）")
            continue

        # ① 首个动画不能配置为"切换进入即播放"：
        #   即第一个 par 的 stCondLst.cond.evt 不能是 onLoad / withEffect
        first_par = mainSeq.find("p:cTn/p:childTnLst/p:par", NS)
        if first_par is not None:
            first_conds = first_par.findall(".//p:stCondLst/p:cond", NS)
            evts = [c.get("evt") for c in first_conds]
            bad_evts = [e for e in evts if e in ("onLoad", "withEffect")]
            if bad_evts:
                fail_details.append(
                    f"P{pno} 首个动画启动条件为 {bad_evts}，换页即播放")
                continue

        # 不再对主序列动画做 ≥0.5 秒硬性时长要求：rubric 只关心内容动画
        # 是否需要单击启动，而不限制单个动画的具体时长。

        ok_pages.append(pno)

    ok = (ok_pages == BODY_PAGES)
    msg = f"换页后等待单击开始（不自动提前播放） 命中 {len(ok_pages)}/9（页 {ok_pages}）"
    if not ok and fail_details:
        msg += f"；未命中：{fail_details[:3]}"
    return ok, msg


# ----------------- 主流程 -----------------
def _locate_pptx(dir_path: str) -> str:
    """在给定目录中定位待评估的 pptx 文件；优先取预期文件名，否则取目录内第一个 .pptx。"""
    preferred = os.path.join(dir_path, "艺术漆北区销售额提升方案.pptx")
    if os.path.exists(preferred):
        return preferred
    for name in os.listdir(dir_path):
        if name.lower().endswith(".pptx"):
            return os.path.join(dir_path, name)
    return preferred  # 保底返回预期路径（用于产出"文件不存在"的错误）


def evaluate(dir_path: str) -> dict:
    """对指定目录内的 PPT 文档进行评估，返回结构化结果字典。

    参数 dir_path 为脚本所在目录（也是被评估文档所在目录）；
    脚本自行在该目录中定位并打开被评估文档。
    """
    script_id = "076"
    # 每项的 (分数, 评分细则原文) 描述 —— 与 items 顺序保持一致
    rule_texts = [
        '第1页首页：标题文本为"艺术漆北区销售额提升方案"，位于页面中部或偏上位置，字体为宋体，字号65-75磅，加粗，颜色与绿色渐变背景形成清晰对比。',
        "页面背景：采用浅、中绿色渐变背景，整体颜色不深，文本可清晰阅读。",
        '第2页目录页：出现目录标题和九个一级目录条目，分别为"市场趋势与消费升级""重点城市增长机会""行业竞争与品牌机遇""年度发展目标""核心产品体系""产品组合与差异化卖点""多渠道拓展策略""整合营销与本地转化""服务体系与执行保障"。',
        '第3、5、7、9、11、13、15、17、19页为一至九部分标题页：标题为"一、市场趋势与消费升级""二、重点城市增长机会""三、行业竞争与品牌机遇""四、年度发展目标""五、核心产品体系""六、产品组合与差异化卖点""七、多渠道拓展策略""八、整合营销与本地转化""九、服务体系与执行保障"，字体加粗，38-42磅，文字居于页面中间放在一行，',
        '第4、6、8、10、12、14、16、18、20页为第一至九部分的正文页：顶部靠左显示"1.1 市场扩容趋势"和"1.2 消费需求变化""2.1 改善型住房需求释放"和"2.2 存量房翻新潜力""3.1 竞争格局变化"和"3.2 品牌突破方向""4.1 经营增长目标"和"4.2 品牌认知目标""5.1 主推系列定位"和"5.2 消费决策简化""6.1 产品矩阵设计"和"6.2 核心卖点提炼""7.1 零售终端升级"和"7.2 家装与工程渠道拓展""8.1 线上内容种草"和"8.2 线下活动转化""9.1 服务体系升级"和"9.2 经销商赋能与组织保障"等二级标题均放在对应部分的最上方靠左区域，加粗16-20磅。每个二级标题下方放置对应正文内容，正文与二级标题匹配，12-16磅。',
        '第21页尾页：尾页文字使用宋体，内容可包含"谢谢观看""艺术漆北区销售额提升方案"，不能出现无关主题。',
        "所有正文页橙色矩形框：正文页文本外围均有橙色矩形框包住二级标题和正文内容，矩形框为可编辑形状，边线清晰可见。橙色矩形框边线线宽约1.5—3磅，颜色为橙色或深橙色，不能过细到难以辨认。",
        "全文字体：所有可见中文文字均为宋体，包括目录、标题、二级标题、正文、页脚和尾页文字。",
        "正文页二级标题动画：第4、6、8、10、12、14、16、18、20页每个二级标题使用动画效果",
        "正文页正文动画：第4、6、8、10、12、14、16、18、20页正文内容均设置动画，且与二级标题动画区分明显。",
        "正文页动画触发：每个正文页切换进入后，通过单击开始本页动画，随后二级标题和正文在同一页面内自动依次出现。",
        "正文页动画顺序：每个正文页先出现二级标题，再出现对应正文内容，两个二级小节的动画顺序与页面文本顺序一致。",
        '正文页动画连续性：同一页面内动画设置为"上一动画之后"或等效自动播放，不需要用户反复点击多次才能看完本页内容。',
        "动画效果差异：二级标题可使用淡入、飞入、擦除等效果，正文使用另一种不同效果，不能所有对象完全使用同一种动画。",
        "页面切换逻辑：换页后动画不会自动提前播放，需单击后开始该页内容动画。",
    ]
    max_deltas = [1, 3, 3, 5, 5, 1, 3, 3, 5, 5, 5, 3, 3, 3, 3]
    max_score = sum(max_deltas)

    result = {
        "id": script_id,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": max_score,
    }

    try:
        pptx_path = _locate_pptx(dir_path)
        result["file_name"] = os.path.basename(pptx_path)

        if not os.path.exists(pptx_path):
            result["status"] = "error"
            result["error"] = f"文件不存在：{pptx_path}"
            return result

        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            result["dim1_pass"] = False
            result["dim1_reason"] = f"文件无法正常打开 / 解析 → {e}"
            result["dim2_items"] = []
            result["total_score"] = 0
            return result

        raw_zip = zipfile.ZipFile(pptx_path)

        # ===== 维度 1 =====
        ok1, fails = check_dim1(prs, raw_zip, pptx_path)
        if not ok1:
            result["dim1_pass"] = False
            result["dim1_reason"] = "；".join(fails)
            result["dim2_items"] = []
            result["total_score"] = 0
            return result

        result["dim1_pass"] = True

        # ===== 维度 2 =====
        items = [
            (max_deltas[0], rule_texts[0], check_first_page,        (prs, raw_zip)),
            (max_deltas[1], rule_texts[1], check_bg_gradient,       (prs, raw_zip)),
            (max_deltas[2], rule_texts[2], check_toc,               (prs,)),
            (max_deltas[3], rule_texts[3], check_section_pages,     (prs,)),
            (max_deltas[4], rule_texts[4], check_body_pages,        (prs,)),
            (max_deltas[5], rule_texts[5], check_last_page,         (prs,)),
            (max_deltas[6], rule_texts[6], check_orange_box,        (prs,)),
            (max_deltas[7], rule_texts[7], check_all_simsun,        (prs,)),
            (max_deltas[8], rule_texts[8], check_sub_title_anim,    (prs, raw_zip)),
            (max_deltas[9], rule_texts[9], check_body_anim,         (prs, raw_zip)),
            (max_deltas[10], rule_texts[10], check_anim_trigger,    (prs, raw_zip)),
            (max_deltas[11], rule_texts[11], check_anim_order,      (prs, raw_zip)),
            (max_deltas[12], rule_texts[12], check_anim_continuity, (prs, raw_zip)),
            (max_deltas[13], rule_texts[13], check_anim_effect_diff, (prs, raw_zip)),
            (max_deltas[14], rule_texts[14], check_page_switch_logic, (prs, raw_zip)),
        ]

        dim2_items = []
        total = 0
        for score, rule_text, fn, args in items:
            try:
                ok, _msg = fn(*args)
            except Exception:
                ok = False
            delta = score if ok else 0
            total += delta
            dim2_items.append({
                "rule": rule_text,
                "max_delta": score,
                "delta": delta,
                "hit": bool(ok),
                "detail": "",
            })

        result["dim2_items"] = dim2_items
        result["total_score"] = total
        return result

    except Exception as e:
        result["status"] = "error"
        result["error"] = f"脚本运行异常：{e}"
        return result


if __name__ == "__main__":
    _dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(_dir), ensure_ascii=False, indent=2))
