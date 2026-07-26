# -*- coding: utf-8 -*-
"""
officeval_072 自动评估脚本 (PPT 打分)。

统一接口:
    evaluate(dir_path: str) -> dict
接收"脚本所在目录的路径", 由脚本自身在该目录下定位并打开被评估文档,
返回结构化字典 (详见 "脚本接口差异与统一建议.md §2.2")。
"""

import os
import re
import sys
import json
import zipfile

from pptx import Presentation
from lxml import etree

SCRIPT_ID = "072"

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
NS = {'a': A_NS, 'p': P_NS, 'r': R_NS}

# ---------------------------------------------------------------------------
# 工具函数
# ---------------------------------------------------------------------------

def get_run_info(r_elem, p_elem):
    """提取一个 run 的有效 字号(pt)、字体名、加粗 (沿 run -> paragraph defRPr 继承)。"""
    sz = None
    font_name = None
    bold = None

    rPr = r_elem.find('a:rPr', NS)
    if rPr is not None:
        sz = rPr.get('sz')
        b = rPr.get('b')
        if b is not None:
            bold = (b == '1')
        for tag in ('a:latin', 'a:ea', 'a:cs'):
            f = rPr.find(tag, NS)
            if f is not None and f.get('typeface'):
                font_name = f.get('typeface')
                break

    if sz is None or font_name is None or bold is None:
        pPr = p_elem.find('a:pPr', NS)
        if pPr is not None:
            defRPr = pPr.find('a:defRPr', NS)
            if defRPr is not None:
                if sz is None:
                    sz = defRPr.get('sz')
                if bold is None:
                    b = defRPr.get('b')
                    if b is not None:
                        bold = (b == '1')
                if font_name is None:
                    for tag in ('a:latin', 'a:ea', 'a:cs'):
                        f = defRPr.find(tag, NS)
                        if f is not None and f.get('typeface'):
                            font_name = f.get('typeface')
                            break

    sz_pt = (int(sz) / 100.0) if sz else None
    return sz_pt, font_name, bold


def iter_text_runs(slide):
    """遍历 slide 内所有文本 run，yield (text, size_pt, font_name, bold, shape)。"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txBody = shape.text_frame._txBody
        for p in txBody.findall('a:p', NS):
            for r in p.findall('a:r', NS):
                t = r.find('a:t', NS)
                txt = (t.text or '') if t is not None else ''
                if txt.strip():
                    sz, fn, bold = get_run_info(r, p)
                    yield txt, sz, fn, bold, shape


def collect_runs_by_slide(prs):
    """返回 dict: slide_index(1-based) -> list of (text, sz, font, bold, shape)。"""
    runs = {}
    for i, slide in enumerate(prs.slides, 1):
        runs[i] = list(iter_text_runs(slide))
    return runs


def find_runs(runs, keywords, exact=False):
    """从 runs 列表中找出包含给定关键字 (任一) 的 run。"""
    found = []
    for r in runs:
        txt = r[0]
        for kw in keywords:
            if exact:
                if txt.strip() == kw:
                    found.append(r)
                    break
            else:
                if kw in txt:
                    found.append(r)
                    break
    return found


def in_range(sz, lo, hi):
    return sz is not None and lo <= sz <= hi


def check_all_in_range(found_runs, lo, hi):
    """所有命中的 run 字号必须都在 [lo, hi]。空列表视为未命中 -> False。"""
    if not found_runs:
        return False
    return all(in_range(r[1], lo, hi) for r in found_runs)


# ---------------------------------------------------------------------------
# 主评估流程
# ---------------------------------------------------------------------------

# 各条得分/扣分项对应的"评分细则原文"
RULE_TEXTS = {
    "整份 PPT 中字体皆为 微软雅黑 或 等线加粗":
        "整份PPT中字体皆为微软雅黑或等线加粗字体",
    "第2-11页顶部标题左侧编号字体皆为 14-16 磅":
        "第2至11页中顶部标题左侧的编号“01”“02”等字体皆为14-16磅",
    "第2-11页顶部标题字体皆为 22-26 磅":
        "第2至11页中顶部标题“班级初印象：安静里藏着光”等字体皆为22-26磅",
    "第2-11页顶部标题栏右侧小字字体皆为 12-16 磅":
        "第2至11页中顶部标题栏右侧的小字“从观察到行动”等字体皆为12-16磅",
    "第1页封面字体规范":
        "第1页封面：“以光为航·以心育人”字体为34-38磅“晨河班成长共创方案”字体为18-22磅"
        + "“班级文化与学生发展路径”字体为22-26磅“让安静的努力被看见，让微小的改变被珍惜”字体为14-18磅。",
    "第2页字体规范":
        "第2页：包含“观察印象”字体为18-22磅“现状扫描”“带班信念”“核心目标”字体为16-20磅，"
        + "以及标题下方对应的文本内容字体为14-18磅等卡片内容。",
    "第3页字体规范":
        "第3页：包含“待唤醒的星火”“最坚实的土壤”字体为20-24磅，及两个标题下方相应的两大板块"
        + "文本内容字体为12-16磅及底部破局关键文字字体为14-18磅。",
    "第4页字体规范":
        "第4页：中心圆写有“班训 微光汇聚 向阳生长”，“班训”字体为20-24磅“微光汇聚 向阳生长”"
        + "字体为22-26磅，四周保留“发现微光、及时回应、伙伴互助、持续生长”字体皆为14-18磅，"
        + "底部“带班口号”一行字体为14-18磅。",
    "第5页字体规范":
        "第5页：包含“起步阶段”“连接阶段”“跃迁阶段”字体为12-16磅“看见微光、形成光束、自我点亮”"
        + "三个阶段字体为16-20磅，三个阶段下方文本字体为10-14磅及底部目标说明字体为14-18磅。",
    "第6页字体规范":
        "第6页：包含中间部分“成长的内驱力”其中“成长的”字体为14-18磅“内驱力”字体为20-24磅；"
        + "中间圆环和左侧的责任剧场、能量积分、伙伴联盟、心愿花园四个标题字体为14-18磅，"
        + "下方对应文本为10-14磅和右侧“运行原则”字体为16-20磅，下方对应文本内容字体为12-16磅。",
    "第7页字体规范":
        "第7页：包含“把班级...光点。”字体为14-18磅“岗位认领、微光观察员、自我提醒官”字体为14-18磅"
        + "下方对应文本字体为10-14磅和底部成效期待文字字体为12-16磅。",
    "第8页字体规范":
        "第8页：包含中心“能量循环”其中“能量”字体为16-20磅“循环”字体为22-26磅、"
        + "获得光点、存入光册、兑换体验、复盘升级字体为12-16磅下方对应文本字体为10-14磅，"
        + "底部“规则底色”一行字体为14-18磅等内容。",
    "第9页字体规范":
        "第9页：包含六边形同盟图和右侧“搭建方式、核心机制、温柔约定”字体为14-18磅"
        + "下方三张说明卡对应文本字体为10-14磅，左侧“支持”“记录”“观察”等六个版块字体为14-18磅"
        + "“同盟”字体为18-22磅及底部“合作...照亮”字体为14-18磅。",
    "第10页字体规范":
        "第10页：包含“心愿花园”字体为20-24磅“我敢开口了”“我能坚持了”字体为14-18磅"
        + "下方对应文本字体为10-14磅“升级方式”“情感闭环”字体为16-20磅下方对应文本字体为12-16磅等内容。",
    "第11页字体规范":
        "第11页：包含“自信在发声”“习惯在迁移”“氛围在变暖”字体为18-22磅三张卡片下方对应文本"
        + "字体为12-16磅及底部总结语字体为14-18磅。",
    "第12页结束页字体规范":
        "第12页结束页：保留“教育是一场温柔的点亮”字体为30-34磅"
        + "“这份答案沉默而坚定，正被一颗颗愿意尝试的心慢慢写出来。”字体为18-22磅"
        + "“感谢聆听与同行”字体为14-18磅等文字。",
    "第2-11页右下角页码保留、数字与页对应、字体 10-12 磅":
        "整份PPT页码：第2页至第11页右下角页码保留，数字与页面对应，字体为10-12磅。",
    "页面中橙/青/绿三色圆形图案超出边界":
        "页面中橙色，青色，绿色三色圆形图案超出边界",
}

# 全部评分项 (顺序 = 输出顺序); 正分 = 得分项, 负分 = 扣分项
RULES = [
    (5, "整份 PPT 中字体皆为 微软雅黑 或 等线加粗"),
    (5, "第2-11页顶部标题左侧编号字体皆为 14-16 磅"),
    (5, "第2-11页顶部标题字体皆为 22-26 磅"),
    (5, "第2-11页顶部标题栏右侧小字字体皆为 12-16 磅"),
    (3, "第1页封面字体规范"),
    (3, "第2页字体规范"),
    (3, "第3页字体规范"),
    (3, "第4页字体规范"),
    (3, "第5页字体规范"),
    (3, "第6页字体规范"),
    (3, "第7页字体规范"),
    (3, "第8页字体规范"),
    (3, "第9页字体规范"),
    (3, "第10页字体规范"),
    (3, "第11页字体规范"),
    (3, "第12页结束页字体规范"),
    (3, "第2-11页右下角页码保留、数字与页对应、字体 10-12 磅"),
    (-5, "页面中橙/青/绿三色圆形图案超出边界"),
]


def _locate_pptx(dir_path):
    """在 dir_path 目录下定位被评估的 .pptx 文件, 返回绝对路径; 未找到返回 None。"""
    if not os.path.isdir(dir_path):
        return None
    candidates = []
    for name in os.listdir(dir_path):
        low = name.lower()
        if not low.endswith('.pptx'):
            continue
        # 忽略 Office 打开时生成的临时锁文件 (~$ 开头)
        if name.startswith('~$'):
            continue
        candidates.append(os.path.join(dir_path, name))
    if not candidates:
        return None
    # 若有多个, 优先选择修改时间最新的
    candidates.sort(key=os.path.getmtime, reverse=True)
    return candidates[0]


def _evaluate_pptx(path):
    """核心评分流程; 返回 (dim1_pass: bool, dim1_reason: str, hits: list[(delta, name, detail)])。
    hits 仅包含实际命中/触发的项 (正分表示得分项命中, 负分表示扣分项被触发)。
    """
    report_lines = []
    def log(msg):
        # 诊断信息仅收集, 不直接打印; 最终由 evaluate() 汇总
        report_lines.append(msg)

    log("=" * 70)
    log(f"评估文件: {path}")
    log("=" * 70)

    # -----------------------------------------------------------------------
    # 维度 1: 交付文件为 .pptx 格式，文件可正常打开 (一票否决)
    # -----------------------------------------------------------------------
    log("\n【维度 1：交付文件为 .pptx 格式，文件可正常打开】")
    dim1_passed = True
    dim1_reasons = []

    # 1) 文件格式 & 可正常打开
    if not path.lower().endswith('.pptx'):
        dim1_passed = False
        dim1_reasons.append("文件扩展名不是 .pptx")
    if not os.path.exists(path):
        dim1_passed = False
        dim1_reasons.append("文件不存在")

    prs = None
    try:
        prs = Presentation(path)
        log("  ✓ 文件格式为 .pptx，可正常打开")
    except Exception as e:
        dim1_passed = False
        dim1_reasons.append(f"文件无法正常打开: {e}")

    if prs is None:
        log("\n维度 1 不通过，最终得分: 0")
        return False, "; ".join(dim1_reasons) or "无法打开文件", []

    # 页面尺寸 (供后续扣分项判断使用)
    sw, sh = prs.slide_width, prs.slide_height
    log(f"  · 页面尺寸: {sw/914400:.2f} x {sh/914400:.2f} 英寸")

    if not dim1_passed:
        log("\n维度 1 不通过：")
        for r in dim1_reasons:
            log(f"  ✗ {r}")
        log("最终得分: 0")
        return False, "; ".join(dim1_reasons), []

    log("  >> 维度 1 通过\n")

    # -----------------------------------------------------------------------
    # 维度 2: 完成度评分细则
    # -----------------------------------------------------------------------
    log("【维度 2：完成度评分细则】")
    runs_by_slide = collect_runs_by_slide(prs)
    score = 0
    hits = []  # 命中的得分/扣分项

    # 工具：累计加分
    def award(points, name, detail=""):
        nonlocal score
        score += points
        hits.append((points, name, detail))
        sign = "+" if points >= 0 else ""
        log(f"  [{sign}{points}] {name}  {detail}")

    # ---- +5: 整份 PPT 中字体皆为 "微软雅黑" 或 "等线加粗" 字体 ----
    chinese_re = re.compile(r'[一-鿿]')
    # 细则字面拆解 (每一点都要踩到):
    #   范围 = 整份 PPT (12 页) 中的所有文本 run
    #   允许字体 = {"微软雅黑"} ∪ {"等线" 且 加粗}
    #     - "微软雅黑" 不要求加粗
    #     - "等线" 必须加粗 (bold=True) 才合格
    #   不在上述集合内的任何字体 (含 SimSun/宋体/Calibri/Times 等) 均视为违例。
    yahei_aliases = {"Microsoft YaHei", "Microsoft YaHei UI", "微软雅黑",
                     "Microsoft YaHei Light", "微软雅黑 Light"}
    dengxian_aliases = {"DengXian", "等线", "DengXian Light", "等线 Light"}

    def font_ok(fn, bold):
        if fn in yahei_aliases:
            return True              # 微软雅黑，加粗无要求
        if fn in dengxian_aliases:
            return bold is True      # 等线必须加粗
        return False

    bad_runs = []        # (page, text, font, bold)
    total_runs = 0
    for i in range(1, 13):
        for r in runs_by_slide[i]:
            total_runs += 1
            fn = r[2] or ""
            if not font_ok(fn, r[3]):
                bad_runs.append((i, r[0][:20], fn, r[3]))

    if not bad_runs:
        award(5, "整份 PPT 中字体皆为 微软雅黑 或 等线加粗",
              f"共 {total_runs} 个 run 全部满足 (微软雅黑无加粗要求；等线必须加粗)")
    else:
        log(f"  [×]  整份 PPT 中字体皆为 微软雅黑 或 等线加粗 "
            f"(违例 {len(bad_runs)} 处)")
        for v in bad_runs[:5]:
            log(f"        - 第{v[0]}页 '{v[1]}' 字体='{v[2]}' 加粗={v[3]}")

    # ---- +5: 第 2 至 11 页中顶部标题左侧的编号"01""02"等字体皆为 14-16 磅 ----
    # 细则字面拆解 (每一点都要踩到):
    #   范围   = 第 2 至 11 页 (共 10 页)
    #   对象   = "顶部标题左侧的编号" —— 即位于该页顶部、且在主标题"左侧"的两位数字编号
    #   形如   = "01" / "02" / … (两位数字串)
    #   字号   = 皆为 14-16 磅
    #   "皆为" = 该范围 10 页全部满足才得 +5
    NUM_RE = re.compile(r'^\d{2}$')
    page_h_in = (prs.slide_height or 6858000) / 914400
    num_ok_pages = []
    num_bad_pages = []
    for i in range(2, 12):
        slide = prs.slides[i - 1]
        # 1) 定位该页"顶部标题" shape: top 区域 (top < 1/4 页高) 中
        #    主中文标题 (字号最大且文本长度 >= 6 的中文 run 所在的 shape)
        top_shapes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                t_in = shape.top / 914400
            except Exception:
                continue
            if t_in > page_h_in / 4:
                continue
            top_shapes.append(shape)

        title_left = None
        title_max_sz = -1.0
        for shape in top_shapes:
            txBody = shape.text_frame._txBody
            for p in txBody.findall('a:p', NS):
                for r in p.findall('a:r', NS):
                    t = r.find('a:t', NS)
                    txt = (t.text or '') if t is not None else ''
                    if not txt.strip():
                        continue
                    if not chinese_re.search(txt):
                        continue
                    if len(txt.strip()) < 6:
                        continue
                    sz, _, _ = get_run_info(r, p)
                    sz_v = sz if sz is not None else 0
                    if sz_v > title_max_sz:
                        title_max_sz = sz_v
                        title_shape = shape
                        try:
                            title_left = shape.left / 914400
                        except Exception:
                            title_left = None

        # 2) 在顶部区域里找"左侧的两位数字编号"shape
        #    判定：text 整体匹配 ^\d{2}$，且 shape.left 在 title_shape.left 之左
        num_runs = []
        for shape in top_shapes:
            full = shape.text_frame.text.strip()
            if not NUM_RE.match(full):
                continue
            try:
                s_left = shape.left / 914400
            except Exception:
                s_left = None
            if title_left is not None and s_left is not None and s_left >= title_left:
                # 不在标题左侧 -> 跳过 (顶部右侧的编号不属于本细则)
                continue
            # 收集该 shape 中所有 run (用于字号检查)
            txBody = shape.text_frame._txBody
            for p in txBody.findall('a:p', NS):
                for r in p.findall('a:r', NS):
                    t = r.find('a:t', NS)
                    txt = (t.text or '') if t is not None else ''
                    if NUM_RE.match(txt.strip()):
                        sz, _, _ = get_run_info(r, p)
                        num_runs.append((txt.strip(), sz))

        if num_runs and all(in_range(sz, 14, 16) for _, sz in num_runs):
            num_ok_pages.append((i, num_runs[0][0], num_runs[0][1]))
        else:
            num_bad_pages.append((i, num_runs if num_runs else None))

    if len(num_ok_pages) == 10:
        award(5, "第2-11页顶部标题左侧编号字体皆为 14-16 磅",
              f"满足页: {[(p, s) for p, _, s in num_ok_pages]}")
    else:
        log(f"  [×]  第2-11页顶部标题左侧编号字体皆为 14-16 磅 "
            f"(不满足: {num_bad_pages})")

    # ---- +5: 第 2 至 11 页中顶部标题"班级初印象：安静里藏着光"等字体皆为 22-26 磅 ----
    # 细则字面拆解 (每一点都要踩到):
    #   范围   = 第 2 至 11 页 (共 10 页)
    #   对象   = "顶部标题" —— 位于该页顶部、且为主中文标题
    #             细则示例 = "班级初印象：安静里藏着光" 等
    #   字号   = 皆为 22-26 磅
    #   "皆为" = 该范围 10 页全部满足才得 +5
    title_ok = True
    title_detail = []
    for i in range(2, 12):
        slide = prs.slides[i - 1]
        # 顶部区域 (top < 页高 / 4) 内的所有可编辑文本 shape
        top_shapes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                t_in = shape.top / 914400
            except Exception:
                continue
            if t_in > page_h_in / 4:
                continue
            top_shapes.append(shape)

        # 在顶部 shape 中找"主中文标题":
        #   候选 = 文本含中文、长度 >= 6 的 run；选字号最大者作为顶部标题
        title_runs = []
        for shape in top_shapes:
            txBody = shape.text_frame._txBody
            for p in txBody.findall('a:p', NS):
                for r in p.findall('a:r', NS):
                    t = r.find('a:t', NS)
                    txt = (t.text or '') if t is not None else ''
                    if not txt.strip():
                        continue
                    if not chinese_re.search(txt):
                        continue
                    if len(txt.strip()) < 6:
                        continue
                    sz, _, _ = get_run_info(r, p)
                    title_runs.append((txt.strip(), sz))

        if not title_runs:
            title_ok = False
            title_detail.append(f"P{i}=未找到顶部标题")
            continue

        # "顶部标题" 取字号最大的那一条 (允许跨多个 run 拼接的标题, 此处取最大代表)
        max_sz = max((sz for _, sz in title_runs if sz is not None), default=None)
        # 同字号的所有 run (常见: 标题被拆成多 run, 它们字号一致)
        same = [(txt, sz) for txt, sz in title_runs if sz == max_sz]
        ok = max_sz is not None and all(in_range(sz, 22, 26) for _, sz in same)
        if ok:
            title_detail.append(f"P{i}='{same[0][0][:10]}'={max_sz}pt")
        else:
            title_ok = False
            title_detail.append(f"P{i}='{same[0][0][:10]}'={max_sz}pt✗")

    if title_ok:
        award(5, "第2-11页顶部标题字体皆为 22-26 磅", "; ".join(title_detail))
    else:
        log(f"  [×]  第2-11页顶部标题字体皆为 22-26 磅 ({'; '.join(title_detail)})")

    # ---- +5: 第 2 至 11 页中顶部标题栏右侧的小字"从观察到行动"等字体皆为 12-16 磅 ----
    # 细则字面拆解 (每一点都要踩到):
    #   范围   = 第 2 至 11 页 (共 10 页)
    #   对象   = "顶部标题栏右侧的小字" —— 位于该页顶部、且在主标题"右侧"的小字
    #             细则示例 = "从观察到行动" 等
    #   字号   = 皆为 12-16 磅
    #   "皆为" = 该范围 10 页全部满足才得 +5
    sub_ok = True
    sub_detail = []
    for i in range(2, 12):
        slide = prs.slides[i - 1]
        # 1) 顶部区域所有可编辑文本 shape (top < 页高 / 4)
        top_shapes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                t_in = shape.top / 914400
            except Exception:
                continue
            if t_in > page_h_in / 4:
                continue
            top_shapes.append(shape)

        # 2) 定位主标题 shape (字号最大且文本长度 >= 6 的中文 run 所在 shape)
        title_shape = None
        title_left = None
        title_max_sz = -1.0
        for shape in top_shapes:
            txBody = shape.text_frame._txBody
            for p in txBody.findall('a:p', NS):
                for r in p.findall('a:r', NS):
                    t = r.find('a:t', NS)
                    txt = (t.text or '') if t is not None else ''
                    if not txt.strip() or not chinese_re.search(txt):
                        continue
                    if len(txt.strip()) < 6:
                        continue
                    sz, _, _ = get_run_info(r, p)
                    sz_v = sz if sz is not None else 0
                    if sz_v > title_max_sz:
                        title_max_sz = sz_v
                        title_shape = shape
                        try:
                            title_left = shape.left / 914400
                        except Exception:
                            title_left = None

        # 3) 收集 "标题栏右侧的小字" —— 顶部、位于主标题右侧、且非主标题本身
        right_runs = []
        for shape in top_shapes:
            if shape is title_shape:
                continue
            try:
                s_left = shape.left / 914400
            except Exception:
                s_left = None
            if title_left is None or s_left is None or s_left <= title_left:
                # 不在主标题右侧 -> 跳过
                continue
            txBody = shape.text_frame._txBody
            for p in txBody.findall('a:p', NS):
                for r in p.findall('a:r', NS):
                    t = r.find('a:t', NS)
                    txt = (t.text or '') if t is not None else ''
                    if not txt.strip() or not chinese_re.search(txt):
                        continue
                    sz, _, _ = get_run_info(r, p)
                    right_runs.append((txt.strip(), sz))

        if right_runs and all(in_range(sz, 12, 16) for _, sz in right_runs):
            sub_detail.append(f"P{i}='{right_runs[0][0][:10]}'={right_runs[0][1]}pt")
        else:
            sub_ok = False
            sub_detail.append(f"P{i}!={right_runs}")

    if sub_ok:
        award(5, "第2-11页顶部标题栏右侧小字字体皆为 12-16 磅", "; ".join(sub_detail))
    else:
        log(f"  [×]  第2-11页顶部标题栏右侧小字字体皆为 12-16 磅 ({'; '.join(sub_detail)})")

    # ---- +3: 第 1 页封面 ----
    # 细则字面拆解 (每一点都要踩到):
    #   "以光为航·以心育人"               字体 34-38 磅
    #   "晨河班成长共创方案"               字体 18-22 磅
    #   "班级文化与学生发展路径"           字体 22-26 磅
    #   "让安静的努力被看见，让微小的改变被珍惜"  字体 14-18 磅
    # "字体为 X-Y 磅" 即 该文字所有 run 的字号都必须在 [X, Y] 区间;
    # 4 项全部满足才得 +3
    p1 = runs_by_slide[1]
    # 包含原文中"·"，文件里可能是全角"·"或半角"·"；用规范化函数兜底
    def _norm(s):
        return re.sub(r'\s+', '', s).replace('·', '·')
    p1_checks = [
        ("以光为航·以心育人", 34, 38),
        ("晨河班成长共创方案", 18, 22),
        ("班级文化与学生发展路径", 22, 26),
        ("让安静的努力被看见，让微小的改变被珍惜", 14, 18),
    ]
    p1_results = []
    p1_ok = True
    for target, lo, hi in p1_checks:
        target_norm = _norm(target)
        # 在第 1 页找文本规范化后等于 / 包含目标的 run
        matched = []
        for r in p1:
            rtext_norm = _norm(r[0])
            if rtext_norm == target_norm or target_norm in rtext_norm or rtext_norm in target_norm:
                # 后两种条件用于覆盖：原文整段在一个 run 中，或被拆成多个 run 而每段是其子串
                if rtext_norm:  # 非空
                    matched.append(r)
        ok = bool(matched) and all(in_range(r[1], lo, hi) for r in matched)
        p1_results.append(f"'{target[:8]}'={[r[1] for r in matched] if matched else None}pt "
                          f"({lo}-{hi}){'✓' if ok else '✗'}")
        if not ok:
            p1_ok = False
    if p1_ok:
        award(3, "第1页封面字体规范", "; ".join(p1_results))
    else:
        log(f"  [×]  第1页封面字体规范 ({'; '.join(p1_results)})")

    # ---- +3: 第 2 页 ----
    # 细则字面拆解 (每一点都要踩到):
    #   "观察印象"                              字体 18-22 磅
    #   "现状扫描"  "带班信念"  "核心目标"      字体 16-20 磅
    #   以上 4 个标题"下方对应的文本内容"        字体 14-18 磅
    #   (全部满足才得 +3)
    slide2 = prs.slides[1]
    p2_results = []
    p2_ok = True

    # 收集第 2 页所有可编辑文本 shape (含位置)
    s2_shapes = []
    for shape in slide2.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        try:
            s2_shapes.append((
                shape,
                shape.left / 914400,
                shape.top / 914400,
                shape.width / 914400,
                shape.height / 914400,
                shape.text_frame.text.strip(),
            ))
        except Exception:
            pass

    def find_shape_by_text(target):
        for s in s2_shapes:
            if s[5].replace('\n', '').startswith(target) or target in s[5].replace('\n', ''):
                return s
        return None

    def shape_run_sizes(shape):
        sizes = []
        txBody = shape.text_frame._txBody
        for p in txBody.findall('a:p', NS):
            for r in p.findall('a:r', NS):
                t = r.find('a:t', NS)
                txt = (t.text or '') if t is not None else ''
                if not txt.strip():
                    continue
                sz, _, _ = get_run_info(r, p)
                sizes.append((txt.strip(), sz))
        return sizes

    def find_body_shape_below(title_shape_info):
        """在标题 shape 下方、且水平区间有交集 (同一卡片) 中找最近的非标题 shape。"""
        if title_shape_info is None:
            return None
        _, tl, tt, tw, th, _ = title_shape_info
        candidates = []
        for s in s2_shapes:
            shape, sl, st_, sw_, sh_, stext = s
            if s is title_shape_info:
                continue
            # 在标题正下方: 垂直在标题之下 & 水平区间相交
            if st_ < tt + th * 0.5:
                continue
            if sl + sw_ < tl or sl > tl + tw:
                continue
            # 距离差
            candidates.append((st_ - tt, s))
        if not candidates:
            return None
        candidates.sort(key=lambda x: x[0])
        return candidates[0][1]

    title_specs = [
        ("观察印象", 18, 22),
        ("现状扫描", 16, 20),
        ("带班信念", 16, 20),
        ("核心目标", 16, 20),
    ]
    for title_text, lo, hi in title_specs:
        info = find_shape_by_text(title_text)
        if info is None:
            p2_ok = False
            p2_results.append(f"'{title_text}'=未找到✗")
            continue
        # 标题字号: 该 shape 内所有 run
        sizes = [sz for _, sz in shape_run_sizes(info[0])]
        ok_title = bool(sizes) and all(in_range(sz, lo, hi) for sz in sizes)
        p2_results.append(f"'{title_text}'={sizes}({lo}-{hi}){'✓' if ok_title else '✗'}")
        if not ok_title:
            p2_ok = False
        # 标题下方对应的文本内容 14-18 磅
        body_info = find_body_shape_below(info)
        if body_info is None:
            p2_ok = False
            p2_results.append(f"'{title_text}'正文=未找到✗")
            continue
        body_sizes = [sz for _, sz in shape_run_sizes(body_info[0])]
        ok_body = bool(body_sizes) and all(in_range(sz, 14, 18) for sz in body_sizes)
        p2_results.append(f"'{title_text}'正文={body_sizes}(14-18){'✓' if ok_body else '✗'}")
        if not ok_body:
            p2_ok = False

    if p2_ok:
        award(3, "第2页字体规范", "; ".join(p2_results))
    else:
        log(f"  [×]  第2页字体规范 ({'; '.join(p2_results)})")

    # ---- +3: 第 3 页 ----
    # ---- +3: 第 3 页 ----
    # 细则字面拆解 (每一点都要踩到):
    #   "待唤醒的星火"  "最坚实的土壤"            字体 20-24 磅
    #   两个标题下方"相应的两大板块文本内容"      字体 12-16 磅
    #     - "待唤醒的星火" 板块正文须包含: 学习、纪律、心理
    #     - "最坚实的土壤" 板块正文须包含: 性格、氛围、潜能
    #   底部"破局关键"文字                        字体 14-18 磅
    #     - 完整文本: "破局关键：把"被动安静"转化为"愿意尝试"，
    #                 把"偶然发光"沉淀为"持续成长""
    #     - 判定关键短语齐备: 破局关键 / 被动安静 / 愿意尝试 / 偶然发光 / 持续成长
    #   (全部满足才得 +3)
    slide3 = prs.slides[2]
    p3_results = []
    p3_ok = True

    s3_shapes = []
    for shape in slide3.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        try:
            s3_shapes.append((
                shape,
                shape.left / 914400,
                shape.top / 914400,
                shape.width / 914400,
                shape.height / 914400,
                shape.text_frame.text.strip(),
            ))
        except Exception:
            pass

    def s3_find(target):
        for s in s3_shapes:
            if target in s[5].replace('\n', ''):
                return s
        return None

    def s3_sizes(shape):
        sizes = []
        txBody = shape.text_frame._txBody
        for p in txBody.findall('a:p', NS):
            for r in p.findall('a:r', NS):
                t = r.find('a:t', NS)
                txt = (t.text or '') if t is not None else ''
                if not txt.strip():
                    continue
                sz, _, _ = get_run_info(r, p)
                sizes.append(sz)
        return sizes

    def s3_bodies_below(title_info):
        """收集位于该标题正下方、与该标题水平区间相交的所有非标题 shape (两大板块的多段正文)。"""
        results = []
        _, tl, tt, tw, th, _ = title_info
        for s in s3_shapes:
            if s is title_info:
                continue
            _, sl, st_, sw_, sh_, _ = s
            if st_ < tt + th * 0.5:        # 必须在标题下方
                continue
            if sl + sw_ < tl or sl > tl + tw:  # 水平区间必须与标题相交
                continue
            results.append(s)
        return results

    # 板块正文预期关键词 (每个词都必须在该板块正文合并文本中出现)
    board_keywords: dict[str, tuple[str, ...]] = {
        "待唤醒的星火": ("学习", "纪律", "心理"),
        "最坚实的土壤": ("性格", "氛围", "潜能"),
    }

    # 两个标题
    for title_text in ("待唤醒的星火", "最坚实的土壤"):
        info = s3_find(title_text)
        if info is None:
            p3_ok = False
            p3_results.append(f"'{title_text}'=未找到✗")
            continue
        title_sizes = s3_sizes(info[0])
        ok_t = bool(title_sizes) and all(in_range(sz, 20, 24) for sz in title_sizes)
        p3_results.append(f"'{title_text}'={title_sizes}(20-24){'✓' if ok_t else '✗'}")
        if not ok_t:
            p3_ok = False
        # 对应板块的文本内容 (可能有多段, 全部检查)
        bodies = s3_bodies_below(info)
        # 排除"破局关键"那段 (它是底部跨整页, 不属于板块下方对应内容)
        bodies = [b for b in bodies if "破局关键" not in b[5]]
        if not bodies:
            p3_ok = False
            p3_results.append(f"'{title_text}'板块正文=未找到✗")
            continue
        all_sizes = []
        merged_parts: list[str] = []
        for b in bodies:
            piece = b[5]  # type: ignore[assignment]
            merged_parts.append(piece if isinstance(piece, str) else "")
            all_sizes.extend(s3_sizes(b[0]))  # type: ignore[arg-type]
        merged_text: str = "".join(merged_parts)
        ok_b = bool(all_sizes) and all(in_range(sz, 12, 16) for sz in all_sizes)
        # 校验板块正文预期内容: 关键词全部齐备
        expected_kws: tuple[str, ...] = board_keywords[title_text]
        merged_clean: str = re.sub(r'\s+', '', merged_text)
        missing_kws: list[str] = [kw for kw in expected_kws if kw not in merged_clean]
        ok_kw = not missing_kws
        kw_mark: str = '✓' if ok_kw else ('缺' + str(missing_kws) + '✗')
        sz_mark: str = '✓' if ok_b else '✗'
        expected_list: list[str] = list(expected_kws)
        line: str = "'{0}'板块={1}(12-16){2} 内容需含{3}{4}".format(
            title_text, all_sizes, sz_mark, expected_list, kw_mark,
        )  # type: ignore[arg-type]
        p3_results.append(line)  # type: ignore[arg-type]
        if not ok_b:
            p3_ok = False
        if not ok_kw:
            p3_ok = False

    # 底部 "破局关键" 文字
    # 完整文本: 破局关键：把"被动安静"转化为"愿意尝试"，把"偶然发光"沉淀为"持续成长"
    # 判定: 该 shape 文本中关键短语必须齐备
    bot_expected = ("破局关键", "被动安静", "愿意尝试", "偶然发光", "持续成长")
    bot_info = s3_find("破局关键")
    if bot_info is None:
        p3_ok = False
        p3_results.append("'破局关键'=未找到✗")
    else:
        bot_sizes = s3_sizes(bot_info[0])
        ok_bot = bool(bot_sizes) and all(in_range(sz, 14, 18) for sz in bot_sizes)
        # 完整文本校验: 五个关键短语必须齐备 (对引号/空白不敏感)
        bot_clean = re.sub(r'\s+', '', bot_info[5])
        missing_phrases = [p for p in bot_expected if p not in bot_clean]
        ok_bot_text = not missing_phrases
        p3_results.append(
            f"'破局关键'={bot_sizes}(14-18){'✓' if ok_bot else '✗'} "
            f"完整文本需含{list(bot_expected)}"
            f"{'✓' if ok_bot_text else f'缺{missing_phrases}✗'}"
        )
        if not ok_bot:
            p3_ok = False
        if not ok_bot_text:
            p3_ok = False

    if p3_ok:
        award(3, "第3页字体规范", "; ".join(p3_results))
    else:
        log(f"  [×]  第3页字体规范 ({'; '.join(p3_results)})")

    # ---- +3: 第 4 页 ----
    # ---- +3: 第 4 页 ----
    # 细则字面拆解 (每一点都要踩到):
    #   中心圆: 写有 "班训 微光汇聚 向阳生长"
    #     - "班训"             字体 20-24 磅
    #     - "微光汇聚 向阳生长" 字体 22-26 磅
    #   四周保留: "发现微光、及时回应、伙伴互助、持续生长" 字体皆为 14-18 磅
    #   底部 "带班口号" 一行     字体 14-18 磅
    #   (全部满足才得 +3)
    slide4 = prs.slides[3]
    p4_results = []
    p4_ok = True

    # 收集第 4 页所有可编辑文本 shape (含位置 + 文本)
    s4_shapes = []
    for shape in slide4.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        try:
            s4_shapes.append((
                shape,
                shape.left / 914400,
                shape.top / 914400,
                shape.width / 914400,
                shape.height / 914400,
                shape.text_frame.text.strip(),
            ))
        except Exception:
            pass

    def s4_sizes(shape):
        sizes = []
        txBody = shape.text_frame._txBody
        for p in txBody.findall('a:p', NS):
            for r in p.findall('a:r', NS):
                t = r.find('a:t', NS)
                txt = (t.text or '') if t is not None else ''
                if not txt.strip():
                    continue
                sz, _, _ = get_run_info(r, p)
                sizes.append((txt.strip(), sz))
        return sizes

    # 1) 找"中心圆": 同一 shape 内同时含 "班训"、"微光汇聚"、"向阳生长"
    center_info = None
    for s in s4_shapes:
        t = s[5].replace('\n', '').replace(' ', '')
        if "班训" in t and "微光汇聚" in t and "向阳生长" in t:
            center_info = s
            break
    if center_info is None:
        p4_ok = False
        p4_results.append("中心圆='班训/微光汇聚/向阳生长'=未找到✗")
    else:
        center_runs = s4_sizes(center_info[0])
        # 按文本内容分两组
        bx_sizes  = [sz for txt, sz in center_runs if "班训" in txt]
        wgxy_sizes = [sz for txt, sz in center_runs
                      if ("微光汇聚" in txt) or ("向阳生长" in txt)]
        ok_bx = bool(bx_sizes) and all(in_range(sz, 20, 24) for sz in bx_sizes)
        ok_wg = bool(wgxy_sizes) and all(in_range(sz, 22, 26) for sz in wgxy_sizes)
        p4_results.append(f"中心圆'班训'={bx_sizes}(20-24){'✓' if ok_bx else '✗'}")
        p4_results.append(f"中心圆'微光汇聚/向阳生长'={wgxy_sizes}(22-26){'✓' if ok_wg else '✗'}")
        if not ok_bx: p4_ok = False
        if not ok_wg: p4_ok = False

    # 2) 四周保留 "发现微光、及时回应、伙伴互助、持续生长" 字体皆为 14-18 磅
    four_words = ("发现微光", "及时回应", "伙伴互助", "持续生长")
    four_results = []
    for w in four_words:
        target = None
        for s in s4_shapes:
            if s is center_info:
                continue
            if w in s[5]:
                target = s
                break
        if target is None:
            p4_ok = False
            four_results.append(f"'{w}'=未找到✗")
            continue
        sizes = [sz for _, sz in s4_sizes(target[0])]
        # 仅判定该 shape 自身 (其文本即四周保留文字)
        ok = bool(sizes) and all(in_range(sz, 14, 18) for sz in sizes)
        four_results.append(f"'{w}'={sizes}{'✓' if ok else '✗'}")
        if not ok: p4_ok = False
    p4_results.append("四周(14-18): " + " ".join(four_results))

    # 3) 底部 "带班口号" 一行 字体 14-18 磅
    bot_info = None
    for s in s4_shapes:
        if "带班口号" in s[5]:
            bot_info = s
            break
    if bot_info is None:
        p4_ok = False
        p4_results.append("'带班口号'=未找到✗")
    else:
        sizes = [sz for _, sz in s4_sizes(bot_info[0])]
        ok = bool(sizes) and all(in_range(sz, 14, 18) for sz in sizes)
        p4_results.append(f"底部'带班口号'={sizes}(14-18){'✓' if ok else '✗'}")
        if not ok: p4_ok = False

    if p4_ok:
        award(3, "第4页字体规范", "; ".join(p4_results))
    else:
        log(f"  [×]  第4页字体规范 ({'; '.join(p4_results)})")

    # ---- +3: 第 5 页 ----
    # 细则字面拆解 (每一点都要踩到):
    #   "起步阶段" "连接阶段" "跃迁阶段"             字体 12-16 磅
    #   "看见微光" "形成光束" "自我点亮" 三个阶段     字体 16-20 磅
    #   三个阶段下方文本 (按预期正文锚定)              字体 10-14 磅
    #     - "看见微光" 下方须含: 鼓励尝试 / 安全感 / 被发现 / 敢表达
    #     - "形成光束" 下方须含: 同伴互助 / 任务协作 / 愿承担 / 会合作
    #     - "自我点亮" 下方须含: 自主规划 / 稳定输出 / 能复盘 / 敢负责
    #   底部目标说明 (按预期文本锚定)                 字体 14-18 磅
    #     - 完整文本: "目标不是让每个孩子立刻耀眼，而是让他们逐渐拥有主动生长的力量"
    #     - 关键短语齐备判定: 目标 / 耀眼 / 主动生长 (兼容标点/空白差异)
    #   (全部满足才得 +3)
    slide5 = prs.slides[4]
    p5_results = []
    p5_ok = True

    s5_shapes = []
    for shape in slide5.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        try:
            s5_shapes.append((
                shape,
                shape.left / 914400,
                shape.top / 914400,
                shape.width / 914400,
                shape.height / 914400,
                shape.text_frame.text.strip(),
            ))
        except Exception:
            pass

    def s5_find(target):
        for s in s5_shapes:
            if target in s[5]:
                return s
        return None

    def s5_sizes(shape):
        sizes = []
        txBody = shape.text_frame._txBody
        for p in txBody.findall('a:p', NS):
            for r in p.findall('a:r', NS):
                t = r.find('a:t', NS)
                txt = (t.text or '') if t is not None else ''
                if not txt.strip():
                    continue
                sz, _, _ = get_run_info(r, p)
                sizes.append(sz)
        return sizes

    def s5_body_below(anchor_info):
        """收集 anchor 下方、水平区间相交、且宽度不超过 anchor 宽度 2 倍的所有 shape
           (该阶段下方文本，排除跨整页的底部目标说明)。"""
        if anchor_info is None:
            return []
        _, al, at_, aw_, ah_, _ = anchor_info
        results = []
        for s in s5_shapes:
            if s is anchor_info:
                continue
            _, sl, st_, sw_, sh_, _ = s
            if st_ < at_ + ah_ * 0.5:
                continue
            if sl + sw_ < al or sl > al + aw_:
                continue
            # 宽度限制：底部跨整页说明的 shape 远宽于阶段卡片，将其排除
            if sw_ > aw_ * 2.5:
                continue
            results.append(s)
        return results

    # (1) "起步阶段" "连接阶段" "跃迁阶段" 字体 12-16 磅
    stage_names = ("起步阶段", "连接阶段", "跃迁阶段")
    stage_name_results = []
    for name in stage_names:
        info = s5_find(name)
        if info is None:
            p5_ok = False
            stage_name_results.append(f"'{name}'=未找到✗")
            continue
        sizes = s5_sizes(info[0])
        ok = bool(sizes) and all(in_range(sz, 12, 16) for sz in sizes)
        stage_name_results.append(f"'{name}'={sizes}{'✓' if ok else '✗'}")
        if not ok: p5_ok = False
    p5_results.append("阶段名(12-16): " + " ".join(stage_name_results))

    # (2) "看见微光" "形成光束" "自我点亮" 字体 16-20 磅
    stage_big_names = ("看见微光", "形成光束", "自我点亮")
    stage_big_results = []
    stage_big_infos = {}
    for name in stage_big_names:
        info = s5_find(name)
        stage_big_infos[name] = info
        if info is None:
            p5_ok = False
            stage_big_results.append(f"'{name}'=未找到✗")
            continue
        sizes = s5_sizes(info[0])
        ok = bool(sizes) and all(in_range(sz, 16, 20) for sz in sizes)
        stage_big_results.append(f"'{name}'={sizes}{'✓' if ok else '✗'}")
        if not ok: p5_ok = False
    p5_results.append("三阶段(16-20): " + " ".join(stage_big_results))

    # (3) 三个阶段下方文本 —— 按预期正文关键词锚定, 字体 10-14 磅
    #     每个阶段合并下方所有正文 shape 文本后, 关键词必须齐备
    stage_expected_kws: dict[str, tuple[str, ...]] = {
        "看见微光": ("鼓励尝试", "安全感", "被发现", "敢表达"),
        "形成光束": ("同伴互助", "任务协作", "愿承担", "会合作"),
        "自我点亮": ("自主规划", "稳定输出", "能复盘", "敢负责"),
    }
    below_results = []
    for name in stage_big_names:
        anchor = stage_big_infos.get(name)
        bodies = s5_body_below(anchor)
        if not bodies:
            p5_ok = False
            below_results.append(f"'{name}'下方=未找到✗")
            continue
        all_sz: list = []
        merged_parts: list[str] = []
        for b in bodies:
            piece = b[5]  # type: ignore[assignment]
            merged_parts.append(piece if isinstance(piece, str) else "")
            all_sz.extend(s5_sizes(b[0]))  # type: ignore[arg-type]
        merged_text: str = "".join(merged_parts)
        ok_sz = bool(all_sz) and all(in_range(sz, 10, 14) for sz in all_sz)
        expected: tuple[str, ...] = stage_expected_kws[name]
        merged_clean: str = re.sub(r'\s+', '', merged_text)
        missing: list[str] = [kw for kw in expected if kw not in merged_clean]
        ok_kw = not missing
        sz_mark: str = '✓' if ok_sz else '✗'
        kw_mark: str = '✓' if ok_kw else ('缺' + str(missing) + '✗')
        below_results.append(
            "'{0}'下方={1}{2} 需含{3}{4}".format(
                name, all_sz, sz_mark, list(expected), kw_mark,
            )
        )
        if not ok_sz: p5_ok = False
        if not ok_kw: p5_ok = False
    p5_results.append("阶段下方(10-14): " + " ".join(below_results))

    # (4) 底部目标说明 —— 按预期文本锚定, 字体 14-18 磅
    #     完整文本: "目标不是让每个孩子立刻耀眼，而是让他们逐渐拥有主动生长的力量"
    #     判定关键短语齐备: 目标 / 耀眼 / 主动生长
    bottom_expected: tuple[str, ...] = ("目标", "耀眼", "主动生长")
    bottom_info = None
    for s in s5_shapes:
        text_clean: str = re.sub(r'\s+', '', s[5] if isinstance(s[5], str) else "")
        if all(kw in text_clean for kw in bottom_expected):
            bottom_info = s
            break
    if bottom_info is None:
        # 兜底: 未按内容匹配到时, 仍保留原"最下方非页码 shape"的兜底定位, 便于诊断
        non_pagenum = [s for s in s5_shapes
                       if not re.fullmatch(r'\d{2}', (s[5] if isinstance(s[5], str) else "").strip())]
        bottom_info = max(non_pagenum, key=lambda s: s[2]) if non_pagenum else None
        p5_ok = False
        p5_results.append(
            "底部目标说明=按预期文本未找到✗ "
            + (f"(兜底最底部='{bottom_info[5][:12]}')" if bottom_info else "")
        )
    else:
        sizes = s5_sizes(bottom_info[0])
        ok = bool(sizes) and all(in_range(sz, 14, 18) for sz in sizes)
        bot_clean: str = re.sub(r'\s+', '', bottom_info[5] if isinstance(bottom_info[5], str) else "")
        missing_b: list[str] = [kw for kw in bottom_expected if kw not in bot_clean]
        ok_text = not missing_b
        sz_mark = '✓' if ok else '✗'
        text_mark: str = '✓' if ok_text else ('缺' + str(missing_b) + '✗')
        p5_results.append(
            "底部目标说明'{0}'={1}(14-18){2} 需含{3}{4}".format(
                (bottom_info[5][:10] if isinstance(bottom_info[5], str) else ""),
                sizes, sz_mark, list(bottom_expected), text_mark,
            )
        )
        if not ok: p5_ok = False
        if not ok_text: p5_ok = False

    if p5_ok:
        award(3, "第5页字体规范", "; ".join(p5_results))
    else:
        log(f"  [×]  第5页字体规范 ({'; '.join(p5_results)})")

    # ---- +3: 第 6 页 ----
    # 细则字面拆解 (每一点都要踩到):
    #   中间部分 "成长的内驱力":
    #     - "成长的"  字体 14-18 磅
    #     - "内驱力"  字体 20-24 磅
    #   "中间圆环和左侧" 的 责任剧场/能量积分/伙伴联盟/心愿花园 四个标题 字体 14-18 磅
    #   (中间圆环 = 围绕"成长的内驱力"环绕分布的 4 个小标签;
    #    左侧     = 列在页面左侧的 4 张卡片标题)
    #   下方对应文本                                              字体 10-14 磅
    #   右侧 "运行原则"                                           字体 16-20 磅
    #   "运行原则" 下方对应文本内容                                字体 12-16 磅
    #   (全部满足才得 +3)
    slide6 = prs.slides[5]
    p6_results = []
    p6_ok = True

    s6_shapes = []
    for shape in slide6.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        try:
            s6_shapes.append((
                shape,
                shape.left / 914400,
                shape.top / 914400,
                shape.width / 914400,
                shape.height / 914400,
                shape.text_frame.text.strip(),
            ))
        except Exception:
            pass

    def s6_sizes_of_run(shape, text_filter=None):
        """返回 shape 中 (text_filter 命中的) 所有 run 字号; text_filter=None 则返回全部。"""
        sizes = []
        txBody = shape.text_frame._txBody
        for p in txBody.findall('a:p', NS):
            for r in p.findall('a:r', NS):
                t = r.find('a:t', NS)
                txt = (t.text or '') if t is not None else ''
                if not txt.strip():
                    continue
                if text_filter is not None and text_filter not in txt:
                    continue
                sz, _, _ = get_run_info(r, p)
                sizes.append(sz)
        return sizes

    def s6_body_below(anchor_info, max_width_factor=2.5):
        if anchor_info is None:
            return []
        _, al, at_, aw_, ah_, _ = anchor_info
        results = []
        for s in s6_shapes:
            if s is anchor_info:
                continue
            _, sl, st_, sw_, sh_, _ = s
            if st_ < at_ + ah_ * 0.5:
                continue
            if sl + sw_ < al or sl > al + aw_:
                continue
            if sw_ > aw_ * max_width_factor:
                continue
            results.append(s)
        return results

    # 1) 中间部分 "成长的内驱力" —— 同一 shape 内含两段
    center_info = None
    for s in s6_shapes:
        t = s[5].replace('\n', '').replace(' ', '')
        if "成长的" in t and "内驱力" in t and len(t) <= 8:
            center_info = s
            break
    if center_info is None:
        p6_ok = False
        p6_results.append("中间'成长的内驱力'=未找到✗")
    else:
        sz_growth = s6_sizes_of_run(center_info[0], "成长的")
        sz_drive  = s6_sizes_of_run(center_info[0], "内驱力")
        ok_a = bool(sz_growth) and all(in_range(sz, 14, 18) for sz in sz_growth)
        ok_b = bool(sz_drive)  and all(in_range(sz, 20, 24) for sz in sz_drive)
        p6_results.append(f"'成长的'={sz_growth}(14-18){'✓' if ok_a else '✗'}")
        p6_results.append(f"'内驱力'={sz_drive}(20-24){'✓' if ok_b else '✗'}")
        if not ok_a: p6_ok = False
        if not ok_b: p6_ok = False

    # 2) "中间圆环 和 左侧" 四个标题 14-18 磅
    # rubric 要求"左侧四标题下方对应文本为 10-14 磅", 因此必须把 圆环标签 与
    # 左侧卡片标题 分开定位: 圆环标签允许没有下方文本, 左侧卡片标题必须存在
    # 对应正文, 缺正文即视为未达标。
    #
    # 分类规则:
    #   左侧卡片标题 = shape 的右边 (left + width) 位于中心圆左侧
    #                 (即整个 shape 完全在中心圆左边; 若中心圆未找到, 退化为 shape
    #                  中心 x < 页宽 * 0.35)
    #   圆环标签    = 其余含有该关键词的 shape (通常紧贴 / 环绕中心圆)
    four_words = ("责任剧场", "能量积分", "伙伴联盟", "心愿花园")
    page_w6 = (prs.slide_width or 12191999) / 914400
    if center_info is not None:
        _, ci_l, _, ci_w, _, _ = center_info
        center_left_edge = ci_l
    else:
        center_left_edge = page_w6 * 0.35

    ring_infos: dict[str, list] = {w: [] for w in four_words}
    left_infos: dict[str, list] = {w: [] for w in four_words}
    for s in s6_shapes:
        if s is center_info:
            continue
        text_clean = s[5].replace('\n', '').replace(' ', '') if isinstance(s[5], str) else ''
        _, sl, _, sw_, _, _ = s
        for w in four_words:
            if w not in text_clean:
                continue
            # 完全在中心圆左边 -> 左侧卡片; 否则 -> 圆环标签
            if sl + sw_ <= center_left_edge:
                left_infos[w].append(s)
            else:
                ring_infos[w].append(s)

    ring_results = []
    for w in four_words:
        infos = ring_infos[w] + left_infos[w]
        if not infos:
            p6_ok = False
            ring_results.append(f"'{w}'=未找到✗")
            continue
        all_sz = []
        for info in infos:
            all_sz.extend(s6_sizes_of_run(info[0]))
        ok = bool(all_sz) and all(in_range(sz, 14, 18) for sz in all_sz)
        ring_results.append(f"'{w}'={all_sz}{'✓' if ok else '✗'}")
        if not ok: p6_ok = False
    p6_results.append("中间圆环和左侧四标题(14-18): " + " ".join(ring_results))

    # 3) 下方对应文本 10-14 磅
    #    - 左侧卡片标题: 必须存在对应正文 (缺失即视为未达标)
    #    - 圆环标签    : 允许没有下方文本 (仅当存在时校验字号)
    body_results = []
    for w in four_words:
        # 左侧卡片: 必须存在正文
        for info in left_infos[w]:
            bodies = s6_body_below(info)
            if not bodies:
                p6_ok = False
                body_results.append(f"左'{w}'下方=未找到✗")
                continue
            all_sz = []
            for b in bodies:
                all_sz.extend(s6_sizes_of_run(b[0]))
            ok = bool(all_sz) and all(in_range(sz, 10, 14) for sz in all_sz)
            body_results.append(f"左'{w}'下方={all_sz}{'✓' if ok else '✗'}")
            if not ok: p6_ok = False
        # 圆环标签: 缺失下方文本不扣分; 存在则校验字号
        for info in ring_infos[w]:
            bodies = s6_body_below(info)
            if not bodies:
                continue
            all_sz = []
            for b in bodies:
                all_sz.extend(s6_sizes_of_run(b[0]))
            ok = bool(all_sz) and all(in_range(sz, 10, 14) for sz in all_sz)
            body_results.append(f"环'{w}'下方={all_sz}{'✓' if ok else '✗'}")
            if not ok: p6_ok = False
    p6_results.append("下方对应文本(10-14): " + " ".join(body_results))

    # 4) 右侧 "运行原则" 16-20 磅
    yxyz_info = None
    for s in s6_shapes:
        if s[5].strip() == "运行原则":
            yxyz_info = s
            break
    if yxyz_info is None:
        p6_ok = False
        p6_results.append("'运行原则'=未找到✗")
    else:
        sz = s6_sizes_of_run(yxyz_info[0])
        ok = bool(sz) and all(in_range(s_, 16, 20) for s_ in sz)
        p6_results.append(f"'运行原则'={sz}(16-20){'✓' if ok else '✗'}")
        if not ok: p6_ok = False

    # 5) "运行原则" 下方对应文本内容 12-16 磅
    if yxyz_info is not None:
        bodies = s6_body_below(yxyz_info, max_width_factor=4)
        if not bodies:
            p6_ok = False
            p6_results.append("'运行原则'下方=未找到✗")
        else:
            all_sz = []
            for b in bodies:
                all_sz.extend(s6_sizes_of_run(b[0]))
            ok = bool(all_sz) and all(in_range(sz, 12, 16) for sz in all_sz)
            p6_results.append(f"'运行原则'下方={all_sz}(12-16){'✓' if ok else '✗'}")
            if not ok: p6_ok = False

    if p6_ok:
        award(3, "第6页字体规范", "; ".join(p6_results))
    else:
        log(f"  [×]  第6页字体规范 ({'; '.join(p6_results)})")

    # ---- +3: 第 7 页 ----
    # 细则字面拆解 (每一点都要踩到):
    #   "把班级...光点。" (顶部一句话)                             字体 14-18 磅
    #   "岗位认领"、"微光观察员"、"自我提醒官" 三个标题             字体 14-18 磅
    #   三个标题下方对应文本                                       字体 10-14 磅
    #   底部 "成效期待" 文字                                        字体 12-16 磅
    #   (全部满足才得 +3)
    slide7 = prs.slides[6]
    p7_results = []
    p7_ok = True

    s7_shapes = []
    for shape in slide7.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        try:
            s7_shapes.append((
                shape,
                shape.left / 914400,
                shape.top / 914400,
                shape.width / 914400,
                shape.height / 914400,
                shape.text_frame.text.strip(),
            ))
        except Exception:
            pass

    def s7_sizes(shape):
        sizes = []
        txBody = shape.text_frame._txBody
        for p in txBody.findall('a:p', NS):
            for r in p.findall('a:r', NS):
                t = r.find('a:t', NS)
                txt = (t.text or '') if t is not None else ''
                if not txt.strip():
                    continue
                sz, _, _ = get_run_info(r, p)
                sizes.append(sz)
        return sizes

    def s7_find(target):
        for s in s7_shapes:
            if target in s[5]:
                return s
        return None

    # (1) 顶部句 "把班级…光点。"
    # 完整预期文本 (rubric 示例上下文):
    #   "把班级日常拆成若干小岗位：每一次认真完成，都会成为被记录的光点。"
    # 判定关键短语齐备: 把班级 / 日常 / 小岗位 / 每一次 / 光点
    top_expected: tuple[str, ...] = ("把班级", "日常", "小岗位", "每一次", "光点")
    top_info = None
    for s in s7_shapes:
        text_clean = re.sub(r'\s+', '', s[5]) if isinstance(s[5], str) else ''
        if all(kw in text_clean for kw in top_expected):
            top_info = s
            break
    if top_info is None:
        p7_ok = False
        p7_results.append(
            "'把班级…光点。'=按预期文本未找到✗ 需含" + str(list(top_expected))
        )
    else:
        sz = s7_sizes(top_info[0])
        ok = bool(sz) and all(in_range(s_, 14, 18) for s_ in sz)
        top_clean: str = re.sub(r'\s+', '', top_info[5] if isinstance(top_info[5], str) else "")
        missing_top: list[str] = [kw for kw in top_expected if kw not in top_clean]
        ok_text = not missing_top
        sz_mark: str = '✓' if ok else '✗'
        text_mark: str = '✓' if ok_text else ('缺' + str(missing_top) + '✗')
        p7_results.append(
            "'把班级…光点。'={0}(14-18){1} 需含{2}{3}".format(
                sz, sz_mark, list(top_expected), text_mark,
            )
        )
        if not ok: p7_ok = False
        if not ok_text: p7_ok = False

    # (2) 三个标题 14-18 磅
    title_words = ("岗位认领", "微光观察员", "自我提醒官")
    title_infos = {}
    title_results = []
    for w in title_words:
        info = s7_find(w)
        title_infos[w] = info
        if info is None:
            p7_ok = False
            title_results.append(f"'{w}'=未找到✗")
            continue
        sz = s7_sizes(info[0])
        ok = bool(sz) and all(in_range(s_, 14, 18) for s_ in sz)
        title_results.append(f"'{w}'={sz}{'✓' if ok else '✗'}")
        if not ok: p7_ok = False
    p7_results.append("三标题(14-18): " + " ".join(title_results))

    # (3) 三标题下方对应文本 10-14 磅 —— 按预期正文关键词锚定
    # 完整预期正文 (rubric 示例):
    #   岗位认领   -> "公开岗位清单，学生按兴趣和能力选择；教师提供轻量指导。"
    #   微光观察员 -> "记录同伴的认真瞬间，传递正向能量，避免只看结果。"
    #   自我提醒官 -> "负责小组节奏与物品整理，让规则变成互相支持。"
    # 关键短语齐备判定 (兼容标点/空白差异):
    body_expected: dict[str, tuple[str, ...]] = {
        "岗位认领":   ("公开岗位清单", "兴趣", "能力", "教师", "轻量指导"),
        "微光观察员": ("记录同伴", "认真瞬间", "正向能量", "只看结果"),
        "自我提醒官": ("小组节奏", "物品整理", "规则", "互相支持"),
    }
    body_results = []
    for w in title_words:
        expected: tuple[str, ...] = body_expected[w]
        # 优先按预期正文关键词全命中来锚定 shape
        body_info = None
        for s in s7_shapes:
            if s is title_infos.get(w):
                continue
            text_clean = re.sub(r'\s+', '', s[5]) if isinstance(s[5], str) else ''
            if all(kw in text_clean for kw in expected):
                body_info = s
                break
        if body_info is None:
            p7_ok = False
            body_results.append(
                f"'{w}'下方=按预期正文未找到✗ 需含{list(expected)}"
            )
            continue
        sz = s7_sizes(body_info[0])
        ok = bool(sz) and all(in_range(s_, 10, 14) for s_ in sz)
        body_results.append(f"'{w}'下方={sz}{'✓' if ok else '✗'} 需含{list(expected)}✓")
        if not ok: p7_ok = False
    p7_results.append("三正文(10-14): " + " ".join(body_results))

    # (4) 底部 "成效期待" 12-16 磅
    bot_info = s7_find("成效期待")
    if bot_info is None:
        p7_ok = False
        p7_results.append("'成效期待'=未找到✗")
    else:
        sz = s7_sizes(bot_info[0])
        ok = bool(sz) and all(in_range(s_, 12, 16) for s_ in sz)
        p7_results.append(f"底部'成效期待'={sz}(12-16){'✓' if ok else '✗'}")
        if not ok: p7_ok = False

    if p7_ok:
        award(3, "第7页字体规范", "; ".join(p7_results))
    else:
        log(f"  [×]  第7页字体规范 ({'; '.join(p7_results)})")

    # ---- +3: 第 8 页 ----
    # 细则字面拆解 (每一点都要踩到):
    #   中心 "能量循环":
    #     - "能量"  字体 16-20 磅
    #     - "循环"  字体 22-26 磅
    #   获得光点、存入光册、兑换体验、复盘升级 (四标题)  字体 12-16 磅
    #   四个标题下方对应文本                              字体 10-14 磅
    #   底部 "规则底色" 一行                              字体 14-18 磅
    #   (全部满足才得 +3)
    slide8 = prs.slides[7]
    p8_results = []
    p8_ok = True

    s8_shapes = []
    for shape in slide8.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        try:
            s8_shapes.append((
                shape,
                shape.left / 914400,
                shape.top / 914400,
                shape.width / 914400,
                shape.height / 914400,
                shape.text_frame.text.strip(),
            ))
        except Exception:
            pass

    def s8_sizes(shape, text_filter=None):
        sizes = []
        txBody = shape.text_frame._txBody
        for p in txBody.findall('a:p', NS):
            for r in p.findall('a:r', NS):
                t = r.find('a:t', NS)
                txt = (t.text or '') if t is not None else ''
                if not txt.strip():
                    continue
                if text_filter is not None and text_filter not in txt:
                    continue
                sz, _, _ = get_run_info(r, p)
                sizes.append(sz)
        return sizes

    def s8_find(target):
        for s in s8_shapes:
            if target in s[5]:
                return s
        return None

    # 1) 中心 "能量循环" —— 按 文本内容 + 布局位于页面中心 双重锚定
    #    文本条件: 同一 shape 中同时含 "能量" 与 "循环"; 使用 run-level 内容检查
    #             而非 len<=6, 避免分行/带空格实现被误伤。
    #    布局条件: 该 shape 的中心点位于页面横向中心 ± 25% 宽度、纵向中心 ± 30% 高度。
    page_w8 = (prs.slide_width or 12191999) / 914400
    page_h8 = (prs.slide_height or 6858000) / 914400
    center_info = None
    for s in s8_shapes:
        text_val = s[5] if isinstance(s[5], str) else ''
        # 通过 run 级别文本判断, 支持分行/空格拆分
        run_texts: list[str] = []
        try:
            txBody = s[0].text_frame._txBody
            for p in txBody.findall('a:p', NS):
                for r in p.findall('a:r', NS):
                    t = r.find('a:t', NS)
                    tt = (t.text or '') if t is not None else ''
                    if tt.strip():
                        run_texts.append(tt)
        except Exception:
            pass
        joined = re.sub(r'\s+', '', "".join(run_texts) or text_val)
        if "能量" not in joined or "循环" not in joined:
            continue
        # 排除该 shape 内还含大量其他非"能量/循环"字符 (避免命中长段正文)
        # 允许标点等少量字符; 严格 <= 8 个非空字符
        if len(joined) > 8:
            continue
        _, sl, st_, sw_, sh_, _ = s
        cx = sl + sw_ / 2
        cy = st_ + sh_ / 2
        if abs(cx - page_w8 / 2) > page_w8 * 0.25:
            continue
        if abs(cy - page_h8 / 2) > page_h8 * 0.30:
            continue
        center_info = s
        break
    if center_info is None:
        p8_ok = False
        p8_results.append("中心'能量循环'=未找到✗ (需 shape 同含'能量'+'循环' 且位于页面中心)")
    else:
        sz_e = s8_sizes(center_info[0], "能量")
        sz_l = s8_sizes(center_info[0], "循环")
        ok_e = bool(sz_e) and all(in_range(sz, 16, 20) for sz in sz_e)
        ok_l = bool(sz_l) and all(in_range(sz, 22, 26) for sz in sz_l)
        p8_results.append(f"中心'能量'={sz_e}(16-20){'✓' if ok_e else '✗'}")
        p8_results.append(f"中心'循环'={sz_l}(22-26){'✓' if ok_l else '✗'}")
        if not ok_e: p8_ok = False
        if not ok_l: p8_ok = False

    # 2) 四个标题 12-16 磅
    four_words = ("获得光点", "存入光册", "兑换体验", "复盘升级")
    four_infos = {}
    four_results = []
    for w in four_words:
        info = s8_find(w)
        four_infos[w] = info
        if info is None:
            p8_ok = False
            four_results.append(f"'{w}'=未找到✗")
            continue
        sz = s8_sizes(info[0])
        ok = bool(sz) and all(in_range(s_, 12, 16) for s_ in sz)
        four_results.append(f"'{w}'={sz}{'✓' if ok else '✗'}")
        if not ok: p8_ok = False
    p8_results.append("四标题(12-16): " + " ".join(four_results))

    # 3) 四标题下方对应文本 10-14 磅 —— 按预期正文关键词锚定
    # 完整预期正文 (rubric 示例):
    #   获得光点 -> "主动表达、完成任务、帮助同伴"
    #   存入光册 -> "以小队为单位记录成长轨迹"
    #   兑换体验 -> "换取展示、选择或合作机会"
    #   复盘升级 -> "把反馈转化为下一次行动"
    # 关键短语齐备判定 (兼容标点/空白差异):
    body_expected8: dict[str, tuple[str, ...]] = {
        "获得光点": ("主动表达", "完成任务", "帮助同伴"),
        "存入光册": ("小队", "记录", "成长轨迹"),
        "兑换体验": ("换取展示", "选择", "合作机会"),
        "复盘升级": ("反馈", "下一次行动"),
    }
    body_results = []
    for w in four_words:
        expected8: tuple[str, ...] = body_expected8[w]
        # 优先在整页中定位同时命中所有关键短语的 shape
        body_info = None
        for s in s8_shapes:
            if s is four_infos.get(w):
                continue
            text_clean = re.sub(r'\s+', '', s[5]) if isinstance(s[5], str) else ''
            if all(kw in text_clean for kw in expected8):
                body_info = s
                break
        if body_info is None:
            p8_ok = False
            body_results.append(
                f"'{w}'下方=按预期正文未找到✗ 需含{list(expected8)}"
            )
            continue
        sz = s8_sizes(body_info[0])
        ok = bool(sz) and all(in_range(s_, 10, 14) for s_ in sz)
        body_results.append(f"'{w}'下方={sz}{'✓' if ok else '✗'} 需含{list(expected8)}✓")
        if not ok: p8_ok = False
    p8_results.append("四正文(10-14): " + " ".join(body_results))

    # 4) 底部 "规则底色" 一行 14-18 磅
    bot_info = s8_find("规则底色")
    if bot_info is None:
        p8_ok = False
        p8_results.append("'规则底色'=未找到✗")
    else:
        sz = s8_sizes(bot_info[0])
        ok = bool(sz) and all(in_range(s_, 14, 18) for s_ in sz)
        p8_results.append(f"底部'规则底色'={sz}(14-18){'✓' if ok else '✗'}")
        if not ok: p8_ok = False

    if p8_ok:
        award(3, "第8页字体规范", "; ".join(p8_results))
    else:
        log(f"  [×]  第8页字体规范 ({'; '.join(p8_results)})")

    # ---- +3: 第 9 页 ----
    # 细则字面拆解 (每一点都要踩到):
    #   右侧 "搭建方式、核心机制、温柔约定"  字体 14-18 磅
    #   右侧三张说明卡下方对应文本           字体 10-14 磅
    #   左侧 "支持""记录""观察" 等 六个版块   字体 14-18 磅
    #   左侧 "同盟"                          字体 18-22 磅
    #   底部 "合作...照亮"                    字体 14-18 磅
    #   (全部满足才得 +3)
    slide9 = prs.slides[8]
    p9_results = []
    p9_ok = True

    s9_shapes = []
    for shape in slide9.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        try:
            s9_shapes.append((
                shape,
                shape.left / 914400,
                shape.top / 914400,
                shape.width / 914400,
                shape.height / 914400,
                shape.text_frame.text.strip(),
            ))
        except Exception:
            pass

    def s9_sizes(shape):
        sizes = []
        txBody = shape.text_frame._txBody
        for p in txBody.findall('a:p', NS):
            for r in p.findall('a:r', NS):
                t = r.find('a:t', NS)
                txt = (t.text or '') if t is not None else ''
                if not txt.strip():
                    continue
                sz, _, _ = get_run_info(r, p)
                sizes.append(sz)
        return sizes

    def s9_find(target):
        for s in s9_shapes:
            if target in s[5]:
                return s
        return None

    def s9_body_below(anchor_info):
        if anchor_info is None:
            return None
        _, al, at_, aw_, ah_, _ = anchor_info
        cands = []
        for s in s9_shapes:
            if s is anchor_info:
                continue
            _, sl, st_, sw_, sh_, _ = s
            if st_ < at_ + ah_ * 0.5:
                continue
            if sl + sw_ < al or sl > al + aw_:
                continue
            cands.append((st_ - at_, s))
        if not cands:
            return None
        cands.sort(key=lambda x: x[0])
        return cands[0][1]

    # 0) 六边形同盟图 (rubric 要求 "包含六边形同盟图")
    #    判定: 第 9 页至少存在 1 个 shape 满足以下任一
    #      - auto_shape_type == MSO_SHAPE.HEXAGON (id 10)
    #      - shape.name (小写) 含 "hex"/"hexagon"/"六边形"
    hexagon_found = False
    hexagon_examples: list[str] = []
    for shape in slide9.shapes:
        # (a) 通过 auto_shape_type 判定
        try:
            ast = getattr(shape, 'auto_shape_type', None)
            if ast is not None and int(ast) == 10:  # MSO_SHAPE.HEXAGON
                hexagon_found = True
                hexagon_examples.append(f"auto_shape=HEXAGON(name={shape.name})")
                break
        except Exception:
            pass
        # (b) 通过名称关键字判定
        try:
            nm = (shape.name or '').lower()
        except Exception:
            nm = ''
        if 'hex' in nm or 'hexagon' in nm or '六边形' in (shape.name or ''):
            hexagon_found = True
            hexagon_examples.append(f"name='{shape.name}'")
            break
    if hexagon_found:
        p9_results.append(f"六边形同盟图✓ ({hexagon_examples[0]})")
    else:
        p9_ok = False
        p9_results.append("六边形同盟图=未找到✗ (需 auto_shape=HEXAGON 或 name 含 hex/hexagon/六边形)")

    # 1) 右侧三标题 14-18 磅
    right_titles = ("搭建方式", "核心机制", "温柔约定")
    right_title_results = []
    right_title_infos = {}
    for w in right_titles:
        info = s9_find(w)
        right_title_infos[w] = info
        if info is None:
            p9_ok = False
            right_title_results.append(f"'{w}'=未找到✗")
            continue
        sz = s9_sizes(info[0])
        ok = bool(sz) and all(in_range(s_, 14, 18) for s_ in sz)
        right_title_results.append(f"'{w}'={sz}{'✓' if ok else '✗'}")
        if not ok: p9_ok = False
    p9_results.append("右侧三标题(14-18): " + " ".join(right_title_results))

    # 2) 右侧三张说明卡对应文本 10-14 磅
    right_body_results = []
    for w in right_titles:
        body_info = s9_body_below(right_title_infos[w])
        if body_info is None:
            p9_ok = False
            right_body_results.append(f"'{w}'下方=未找到✗")
            continue
        sz = s9_sizes(body_info[0])
        ok = bool(sz) and all(in_range(s_, 10, 14) for s_ in sz)
        right_body_results.append(f"'{w}'下方={sz}{'✓' if ok else '✗'}")
        if not ok: p9_ok = False
    p9_results.append("说明卡正文(10-14): " + " ".join(right_body_results))

    # 3) 左侧 "支持""记录""观察" 等 六个版块 14-18 磅
    #    rubric 仅明示 "支持""记录""观察"; 其余三个版块名未见 rubric 明示,
    #    这里改为 —— 必要项: 支持/记录/观察 (三者字号必达标);
    #                任意项: 页面中整体文本恰为单个词的其余 shape 都视作左侧版块
    #                        (若整体存在 6 个及以上的独字/短词 shape 属于同一
    #                         左侧列, 视为通过"六个版块"整体要求)。
    required_words = ("支持", "记录", "观察")
    six_results = []
    # (a) rubric 明示三项: 必须找到且字号达标
    left_short_shapes: list = []
    for w in required_words:
        info = None
        for s in s9_shapes:
            if isinstance(s[5], str) and s[5].strip() == w:
                info = s
                break
        if info is None:
            p9_ok = False
            six_results.append(f"'{w}'=未找到✗")
            continue
        left_short_shapes.append(info)
        sz = s9_sizes(info[0])
        ok = bool(sz) and all(in_range(s_, 14, 18) for s_ in sz)
        six_results.append(f"'{w}'={sz}{'✓' if ok else '✗'}")
        if not ok: p9_ok = False
    # (b) "等" 六个版块: 页面中所有"整体文本为 1-4 字的短词 shape" (排除"同盟" / 页码 / 已算过的三项)
    #     视为其余左侧版块; 数量必须 ≥ 6 (含上面三项) 且它们的字号均达标
    other_left_results = []
    other_left_sizes_ok = True
    for s in s9_shapes:
        if s in left_short_shapes:
            continue
        if not isinstance(s[5], str):
            continue
        tt = s[5].strip()
        if not tt or tt == "同盟":
            continue
        if re.fullmatch(r'\d{1,2}', tt):
            continue
        if len(tt) > 4 or not chinese_re.search(tt):
            continue
        left_short_shapes.append(s)
        sz = s9_sizes(s[0])
        ok = bool(sz) and all(in_range(s_, 14, 18) for s_ in sz)
        other_left_results.append(f"'{tt}'={sz}{'✓' if ok else '✗'}")
        if not ok: other_left_sizes_ok = False
    if len(left_short_shapes) < 6:
        p9_ok = False
        six_results.append(f"其余版块=不足6个 (实到{len(left_short_shapes)})✗")
    elif other_left_results:
        six_results.append("其余版块: " + " ".join(other_left_results))
    if not other_left_sizes_ok:
        p9_ok = False
    p9_results.append("六版块(14-18): " + " ".join(six_results))

    # 4) "同盟" 18-22 磅 (整体文本恰为 "同盟")
    tm_info = None
    for s in s9_shapes:
        if s[5].strip() == "同盟":
            tm_info = s
            break
    if tm_info is None:
        p9_ok = False
        p9_results.append("'同盟'=未找到✗")
    else:
        sz = s9_sizes(tm_info[0])
        ok = bool(sz) and all(in_range(s_, 18, 22) for s_ in sz)
        p9_results.append(f"'同盟'={sz}(18-22){'✓' if ok else '✗'}")
        if not ok: p9_ok = False

    # 5) 底部 "合作...照亮" 14-18 磅
    bot_info = None
    for s in s9_shapes:
        text = s[5].replace('\n', '')
        if text.startswith("合作") and "照亮" in text:
            bot_info = s
            break
    if bot_info is None:
        p9_ok = False
        p9_results.append("底部'合作…照亮'=未找到✗")
    else:
        sz = s9_sizes(bot_info[0])
        ok = bool(sz) and all(in_range(s_, 14, 18) for s_ in sz)
        p9_results.append(f"底部'合作…照亮'={sz}(14-18){'✓' if ok else '✗'}")
        if not ok: p9_ok = False

    if p9_ok:
        award(3, "第9页字体规范", "; ".join(p9_results))
    else:
        log(f"  [×]  第9页字体规范 ({'; '.join(p9_results)})")

    # ---- +3: 第 10 页 ----
    # 细则字面拆解 (每一点都要踩到):
    #   "心愿花园"                              字体 20-24 磅
    #   "我敢开口了" "我能坚持了"                字体 14-18 磅
    #   下方对应文本                             字体 10-14 磅
    #   "升级方式" "情感闭环"                    字体 16-20 磅
    #   下方对应文本                             字体 12-16 磅
    #   (全部满足才得 +3)
    slide10 = prs.slides[9]
    p10_results = []
    p10_ok = True

    s10_shapes = []
    for shape in slide10.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        try:
            s10_shapes.append((
                shape,
                shape.left / 914400,
                shape.top / 914400,
                shape.width / 914400,
                shape.height / 914400,
                shape.text_frame.text.strip(),
            ))
        except Exception:
            pass

    def s10_sizes(shape):
        sizes = []
        txBody = shape.text_frame._txBody
        for p in txBody.findall('a:p', NS):
            for r in p.findall('a:r', NS):
                t = r.find('a:t', NS)
                txt = (t.text or '') if t is not None else ''
                if not txt.strip():
                    continue
                sz, _, _ = get_run_info(r, p)
                sizes.append(sz)
        return sizes

    def s10_find_exact(target):
        for s in s10_shapes:
            if s[5].strip() == target:
                return s
        return None

    # 1) "心愿花园" 20-24 磅 —— 严格匹配整体文本恰为 "心愿花园" 的 shape
    #    (去除脚本假设的"含'心愿花园'且不含'第四步'"兜底)
    xyhy_info = s10_find_exact("心愿花园")
    if xyhy_info is None:
        p10_ok = False
        p10_results.append("'心愿花园'=未找到✗ (需整体文本恰为'心愿花园'的 shape)")
    else:
        sz = s10_sizes(xyhy_info[0])
        ok = bool(sz) and all(in_range(s_, 20, 24) for s_ in sz)
        p10_results.append(f"'心愿花园'={sz}(20-24){'✓' if ok else '✗'}")
        if not ok: p10_ok = False

    # 各卡片下方正文预期内容 (rubric 示例):
    #   我敢开口了 -> "从小声回答到主动补充"
    #   我能坚持了 -> "从提醒完成到自我检查"
    #   升级方式   -> "把阶段性努力做成可展示的成长展：照片换成手绘卡、短语换成学生自己的成长句。"
    #   情感闭环   -> "愿望不以物质奖励为主，更多转化为一次选择权、一次展示机会、一次被认真倾听。"
    body_expected10: dict[str, tuple[str, ...]] = {
        "我敢开口了": ("小声回答", "主动补充"),
        "我能坚持了": ("提醒完成", "自我检查"),
        "升级方式":   ("阶段性努力", "可展示", "成长展", "手绘卡", "成长句"),
        "情感闭环":   ("愿望", "物质奖励", "选择权", "展示机会", "被认真倾听"),
    }

    def _s10_find_body_by_kws(anchor_shape, expected: tuple[str, ...]):
        """在整页 shape 中查找同时包含所有关键短语的 shape (排除锚点自身)。"""
        for s in s10_shapes:
            if s is anchor_shape:
                continue
            text_val = s[5] if isinstance(s[5], str) else ''
            text_clean = re.sub(r'\s+', '', text_val)
            if all(kw in text_clean for kw in expected):
                return s
        return None

    # 2) "我敢开口了" "我能坚持了" 14-18 磅; 下方对应文本按预期正文锚定, 10-14 磅
    cards = ("我敢开口了", "我能坚持了")
    card_results = []
    card_body_results = []
    for w in cards:
        info = s10_find_exact(w)
        if info is None:
            p10_ok = False
            card_results.append(f"'{w}'=未找到✗")
            continue
        sz = s10_sizes(info[0])
        ok = bool(sz) and all(in_range(s_, 14, 18) for s_ in sz)
        card_results.append(f"'{w}'={sz}{'✓' if ok else '✗'}")
        if not ok: p10_ok = False
        # 下方对应文本 —— 按预期正文关键词锚定 (不再取几何最近)
        expected10: tuple[str, ...] = body_expected10[w]
        body_info = _s10_find_body_by_kws(info, expected10)
        if body_info is None:
            p10_ok = False
            card_body_results.append(
                f"'{w}'下方=按预期正文未找到✗ 需含{list(expected10)}"
            )
            continue
        sz_b = s10_sizes(body_info[0])
        ok_b = bool(sz_b) and all(in_range(s_, 10, 14) for s_ in sz_b)
        card_body_results.append(
            f"'{w}'下方={sz_b}{'✓' if ok_b else '✗'} 需含{list(expected10)}✓"
        )
        if not ok_b: p10_ok = False
    p10_results.append("两卡标题(14-18): " + " ".join(card_results))
    p10_results.append("两卡正文(10-14): " + " ".join(card_body_results))

    # 3) "升级方式" "情感闭环" 16-20 磅; 下方对应文本按预期正文锚定, 12-16 磅
    titles2 = ("升级方式", "情感闭环")
    t2_results = []
    t2_body_results = []
    for w in titles2:
        info = s10_find_exact(w)
        if info is None:
            p10_ok = False
            t2_results.append(f"'{w}'=未找到✗")
            continue
        sz = s10_sizes(info[0])
        ok = bool(sz) and all(in_range(s_, 16, 20) for s_ in sz)
        t2_results.append(f"'{w}'={sz}{'✓' if ok else '✗'}")
        if not ok: p10_ok = False
        # 下方对应文本 —— 按预期正文关键词锚定
        expected10 = body_expected10[w]
        body_info = _s10_find_body_by_kws(info, expected10)
        if body_info is None:
            p10_ok = False
            t2_body_results.append(
                f"'{w}'下方=按预期正文未找到✗ 需含{list(expected10)}"
            )
            continue
        sz_b = s10_sizes(body_info[0])
        ok_b = bool(sz_b) and all(in_range(s_, 12, 16) for s_ in sz_b)
        t2_body_results.append(
            f"'{w}'下方={sz_b}{'✓' if ok_b else '✗'} 需含{list(expected10)}✓"
        )
        if not ok_b: p10_ok = False
    p10_results.append("升级方式/情感闭环(16-20): " + " ".join(t2_results))
    p10_results.append("二者正文(12-16): " + " ".join(t2_body_results))

    if p10_ok:
        award(3, "第10页字体规范", "; ".join(p10_results))
    else:
        log(f"  [×]  第10页字体规范 ({'; '.join(p10_results)})")

    # ---- +3: 第 11 页 ----
    # 细则字面拆解 (每一点都要踩到):
    #   "自信在发声" "习惯在迁移" "氛围在变暖" 三张卡片标题  字体 18-22 磅
    #   三张卡片下方对应文本                                字体 12-16 磅
    #   底部总结语                                          字体 14-18 磅
    #   (全部满足才得 +3)
    slide11 = prs.slides[10]
    p11_results = []
    p11_ok = True

    s11_shapes = []
    for shape in slide11.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        try:
            s11_shapes.append((
                shape,
                shape.left / 914400,
                shape.top / 914400,
                shape.width / 914400,
                shape.height / 914400,
                shape.text_frame.text.strip(),
            ))
        except Exception:
            pass

    def s11_sizes(shape):
        sizes = []
        txBody = shape.text_frame._txBody
        for p in txBody.findall('a:p', NS):
            for r in p.findall('a:r', NS):
                t = r.find('a:t', NS)
                txt = (t.text or '') if t is not None else ''
                if not txt.strip():
                    continue
                sz, _, _ = get_run_info(r, p)
                sizes.append(sz)
        return sizes

    def s11_find_exact(target):
        for s in s11_shapes:
            if s[5].strip() == target:
                return s
        return None

    def s11_body_below(anchor_info, max_width_factor=2.0):
        if anchor_info is None:
            return None
        _, al, at_, aw_, ah_, _ = anchor_info
        cands = []
        for s in s11_shapes:
            if s is anchor_info:
                continue
            _, sl, st_, sw_, sh_, _ = s
            if st_ < at_ + ah_ * 0.5:
                continue
            if sl + sw_ < al or sl > al + aw_:
                continue
            if sw_ > aw_ * max_width_factor:
                continue
            cands.append((st_ - at_, s))
        if not cands:
            return None
        cands.sort(key=lambda x: x[0])
        return cands[0][1]

    # 1) 三张卡片标题 18-22 磅
    cards = ("自信在发声", "习惯在迁移", "氛围在变暖")
    card_infos = {}
    card_results = []
    for w in cards:
        info = s11_find_exact(w)
        card_infos[w] = info
        if info is None:
            p11_ok = False
            card_results.append(f"'{w}'=未找到✗")
            continue
        sz = s11_sizes(info[0])
        ok = bool(sz) and all(in_range(s_, 18, 22) for s_ in sz)
        card_results.append(f"'{w}'={sz}{'✓' if ok else '✗'}")
        if not ok: p11_ok = False
    p11_results.append("三卡标题(18-22): " + " ".join(card_results))

    # 2) 三张卡片下方对应文本 12-16 磅 —— 按预期正文关键词锚定 (不再取几何最近下方)
    #    rubric 明示只要求"下方对应文本"字号 12-16, 未直接给出正文原文;
    #    这里按每张卡片标题对应的班级成长意象, 选取语义上必须命中的关键短语,
    #    同时具备"标签性 + 具体化"两类词, 避免与其它卡片文本混淆:
    #      自信在发声 —— 从"沉默不语"到"愿意开口": 发言/回答
    #      习惯在迁移 —— 课堂上的习惯"迁移/延续"到课后与家庭
    #      氛围在变暖 —— 同伴之间"互助/回应/温度"
    #    注: 如后续拿到官方参考文档, 可将下方关键词调整为实际短语。
    body_expected11: dict[str, tuple[str, ...]] = {
        "自信在发声": ("发言", "回答"),
        "习惯在迁移": ("课堂", "课后"),
        "氛围在变暖": ("同伴", "互助"),
    }

    def _s11_find_body_by_kws(anchor_shape, expected: tuple[str, ...]):
        """在整页 shape 中查找同时包含所有关键短语的 shape (排除锚点自身)。"""
        for s in s11_shapes:
            if s is anchor_shape:
                continue
            text_val = s[5] if isinstance(s[5], str) else ''
            text_clean = re.sub(r'\s+', '', text_val)
            if all(kw in text_clean for kw in expected):
                return s
        return None

    body_results = []
    for w in cards:
        anchor = card_infos.get(w)
        expected_body: tuple[str, ...] = body_expected11[w]
        body_info = _s11_find_body_by_kws(anchor, expected_body)
        if body_info is None:
            p11_ok = False
            body_results.append(
                f"'{w}'下方=按预期正文未找到✗ 需含{list(expected_body)}"
            )
            continue
        sz = s11_sizes(body_info[0])
        ok = bool(sz) and all(in_range(s_, 12, 16) for s_ in sz)
        body_results.append(
            f"'{w}'下方={sz}{'✓' if ok else '✗'} 需含{list(expected_body)}✓"
        )
        if not ok: p11_ok = False
    p11_results.append("三卡正文(12-16): " + " ".join(body_results))

    # 3) 底部总结语 14-18 磅 —— 按预期文本锚定 (不再取几何 top 最大兜底)
    #    第 11 页作为"班级文化落地成效"页, 底部总结语意在收束: 从安静到发声、
    #    从个体到集体、微小改变逐渐被看见, 关键短语必须齐备。
    #    注: 关键词依据 rubric 语义拟定; 拿到参考文档后再对齐官方原文。
    bottom_expected11: tuple[str, ...] = ("安静", "发声", "改变", "被看见")
    bottom_info = None
    for s in s11_shapes:
        text_clean = re.sub(r'\s+', '', s[5] if isinstance(s[5], str) else "")
        if all(kw in text_clean for kw in bottom_expected11):
            bottom_info = s
            break
    if bottom_info is None:
        p11_ok = False
        # 兜底: 未按内容匹配到时, 保留原"最下方非页码 shape"用于诊断
        non_pagenum = [
            s for s in s11_shapes
            if not re.fullmatch(r'\d{2}', (s[5] if isinstance(s[5], str) else "").strip())
        ]
        fallback = max(non_pagenum, key=lambda s: s[2]) if non_pagenum else None
        p11_results.append(
            "底部总结语=按预期文本未找到✗ 需含{0}".format(list(bottom_expected11))
            + (f" (兜底最底部='{fallback[5][:12]}')" if fallback else "")
        )
    else:
        sz = s11_sizes(bottom_info[0])
        ok = bool(sz) and all(in_range(s_, 14, 18) for s_ in sz)
        bot_clean = re.sub(r'\s+', '', bottom_info[5] if isinstance(bottom_info[5], str) else "")
        missing_b: list[str] = [kw for kw in bottom_expected11 if kw not in bot_clean]
        ok_text = not missing_b
        sz_mark = '✓' if ok else '✗'
        text_mark = '✓' if ok_text else ('缺' + str(missing_b) + '✗')
        p11_results.append(
            "底部总结语'{0}'={1}(14-18){2} 需含{3}{4}".format(
                (bottom_info[5][:12] if isinstance(bottom_info[5], str) else ""),
                sz, sz_mark, list(bottom_expected11), text_mark,
            )
        )
        if not ok: p11_ok = False
        if not ok_text: p11_ok = False

    if p11_ok:
        award(3, "第11页字体规范", "; ".join(p11_results))
    else:
        log(f"  [×]  第11页字体规范 ({'; '.join(p11_results)})")

    # ---- +3: 第 12 页结束页 ----
    # 细则字面拆解 (每一点都要踩到):
    #   "教育是一场温柔的点亮"                                字体 30-34 磅
    #   "这份答案沉默而坚定，正被一颗颗愿意尝试的心慢慢写出来。"  字体 18-22 磅
    #   "感谢聆听与同行"                                       字体 14-18 磅
    #   (全部满足才得 +3)
    slide12 = prs.slides[11]
    p12_results = []
    p12_ok = True

    s12_shapes = []
    for shape in slide12.shapes:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            continue
        try:
            s12_shapes.append((
                shape,
                shape.left / 914400,
                shape.top / 914400,
                shape.width / 914400,
                shape.height / 914400,
                shape.text_frame.text.strip(),
            ))
        except Exception:
            pass

    def s12_sizes(shape):
        sizes = []
        txBody = shape.text_frame._txBody
        for p in txBody.findall('a:p', NS):
            for r in p.findall('a:r', NS):
                t = r.find('a:t', NS)
                txt = (t.text or '') if t is not None else ''
                if not txt.strip():
                    continue
                sz, _, _ = get_run_info(r, p)
                sizes.append(sz)
        return sizes

    # 规范化: 去空白、忽略换行，便于整段比对 (例如长句被换行分成多 paragraph)
    def s12_norm(s):
        return re.sub(r'\s+', '', s)

    p12_checks = [
        ("教育是一场温柔的点亮", 30, 34),
        ("这份答案沉默而坚定，正被一颗颗愿意尝试的心慢慢写出来。", 18, 22),
        ("感谢聆听与同行", 14, 18),
    ]
    for target, lo, hi in p12_checks:
        target_n = s12_norm(target)
        match = None
        for s in s12_shapes:
            text_n = s12_norm(s[5])
            if text_n == target_n or target_n in text_n or text_n in target_n:
                match = s
                break
        if match is None:
            p12_ok = False
            p12_results.append(f"'{target[:10]}…'=未找到✗")
            continue
        sz = s12_sizes(match[0])
        ok = bool(sz) and all(in_range(s_, lo, hi) for s_ in sz)
        p12_results.append(f"'{target[:10]}…'={sz}({lo}-{hi}){'✓' if ok else '✗'}")
        if not ok:
            p12_ok = False

    if p12_ok:
        award(3, "第12页结束页字体规范", "; ".join(p12_results))
    else:
        log(f"  [×]  第12页结束页字体规范 ({'; '.join(p12_results)})")

    # ---- +3: 整份 PPT 页码: 第 2 页至第 11 页右下角页码保留, 数字与页面对应, 字体 10-12 磅 ----
    # 细则字面拆解 (每一点都要踩到):
    #   范围   = 第 2 页至第 11 页 (共 10 页)
    #   位置   = "右下角" —— shape 位于该页右下区域
    #            (left + width/2 > 页宽/2; top + height/2 > 页高/2)
    #   保留   = 在该右下区域内存在该页码 shape
    #   数字与页面对应 = 该 shape 整体文本严格等于本页页号 (这里按两位数 f"{i:02d}";
    #            也接受单数字 f"{i}", 如第 2 页可写 "2" 或 "02")
    #   字体   = 字号 10-12 磅
    #   (10 页全部满足才得 +3)
    page_w_in = (prs.slide_width or 12191999) / 914400
    page_h_in_2 = (prs.slide_height or 6858000) / 914400
    page_ok_pages = []
    page_bad_pages = []
    for i in range(2, 12):
        slide = prs.slides[i - 1]
        expected_variants = (f"{i:02d}", str(i))
        found = None
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if text not in expected_variants:
                continue
            try:
                l = shape.left / 914400
                t = shape.top / 914400
                w = shape.width / 914400
                h = shape.height / 914400
            except Exception:
                continue
            # "右下角": shape 中心位于页面右下 1/2 ∩ 下 1/2 象限
            cx = l + w / 2
            cy = t + h / 2
            if cx <= page_w_in / 2 or cy <= page_h_in_2 / 2:
                continue
            # 收集字号
            sizes = []
            txBody = shape.text_frame._txBody
            for p in txBody.findall('a:p', NS):
                for r in p.findall('a:r', NS):
                    tt = r.find('a:t', NS)
                    txt = (tt.text or '') if tt is not None else ''
                    if not txt.strip():
                        continue
                    sz, _, _ = get_run_info(r, p)
                    sizes.append(sz)
            found = (text, sizes, (l, t))
            break
        if found is None:
            page_bad_pages.append((i, "未找到右下角页码"))
            continue
        text, sizes, _ = found
        ok_text = text in expected_variants
        ok_size = bool(sizes) and all(in_range(sz, 10, 12) for sz in sizes)
        if ok_text and ok_size:
            page_ok_pages.append((i, text, sizes))
        else:
            page_bad_pages.append((i, f"text={text} sizes={sizes}"))

    if len(page_ok_pages) == 10:
        award(3, "第2-11页右下角页码保留、数字与页对应、字体 10-12 磅",
              f"全部 10 页满足: {[(p, t, s) for p, t, s in page_ok_pages]}")
    else:
        log(f"  [×]  第2-11页右下角页码保留、数字与页对应、字体 10-12 磅 "
            f"(满足: {page_ok_pages}, 不满足: {page_bad_pages})")

    # -----------------------------------------------------------------------
    # 扣分项
    # -----------------------------------------------------------------------
    log("\n  -- 扣分项检查 --")

    # ---- -5: 页面中橙色, 青色, 绿色三色圆形图案超出边界 ----
    # 细则字面拆解 (满足"任意一点"即扣分):
    #   范围 = 页面中的圆形图案 (真几何椭圆/圆形, 含 prstGeom=ellipse 与近似 1:1 自由形状)
    #   颜色 = 橙色 / 青色 / 绿色 三色之一 (支持 srgbClr 与 schemeClr, 显式颜色容差)
    #   超出边界 = shape 任一边 (left/top/right/bottom) 落在页面矩形 [0, 页宽] x [0, 页高] 之外
    sw_in = sw / 914400
    sh_in = sh / 914400

    # ---------- (1) 解析 theme1.xml 的配色方案 (schemeClr -> RGB hex6) ----------
    #   theme 内可能是 <a:srgbClr val="RRGGBB"/> 或 <a:sysClr lastClr="RRGGBB"/>。
    def _read_theme_scheme(pptx_path: str) -> "dict[str, str]":
        scheme: "dict[str, str]" = {}
        try:
            with zipfile.ZipFile(pptx_path) as zf:
                theme_name = None
                for n in zf.namelist():
                    if n.startswith('ppt/theme/theme') and n.endswith('.xml'):
                        theme_name = n
                        break
                if theme_name is None:
                    return scheme
                xml_bytes = zf.read(theme_name)
            root = etree.fromstring(xml_bytes)
            clr_scheme = root.find('.//a:themeElements/a:clrScheme', NS)
            if clr_scheme is None:
                return scheme
            for child in clr_scheme:
                # child.tag 形如 '{...}accent1', 提取本地名作为 key
                local = etree.QName(child.tag).localname
                srgb = child.find('a:srgbClr', NS)
                sys_ = child.find('a:sysClr', NS)
                if srgb is not None and srgb.get('val'):
                    scheme[local] = srgb.get('val').upper()
                elif sys_ is not None and sys_.get('lastClr'):
                    scheme[local] = sys_.get('lastClr').upper()
        except Exception:
            pass
        return scheme

    theme_scheme = _read_theme_scheme(path)

    # ---------- (2) 解析 shape 的填充颜色 (支持 srgbClr / schemeClr + lumMod/lumOff/shade/tint) ----------
    def _apply_lum_mod(hex6: str, clr_elem) -> str:
        """按 lumMod / lumOff / shade / tint 微调, 返回 RRGGBB。粗略实现: 线性 mix。"""
        try:
            r = int(hex6[0:2], 16); g = int(hex6[2:4], 16); b = int(hex6[4:6], 16)
        except Exception:
            return hex6
        # lumMod: 乘系数 (val/100000)
        lm = clr_elem.find('a:lumMod', NS)
        lo = clr_elem.find('a:lumOff', NS)
        sh = clr_elem.find('a:shade', NS)
        ti = clr_elem.find('a:tint', NS)
        def _pct(el):
            try:
                return int(el.get('val')) / 100000.0
            except Exception:
                return None
        if lm is not None and _pct(lm) is not None:
            k = _pct(lm)
            r = int(r * k); g = int(g * k); b = int(b * k)
        if lo is not None and _pct(lo) is not None:
            k = _pct(lo)
            r = int(r + (255 - r) * k)
            g = int(g + (255 - g) * k)
            b = int(b + (255 - b) * k)
        if sh is not None and _pct(sh) is not None:
            k = _pct(sh)
            r = int(r * k); g = int(g * k); b = int(b * k)
        if ti is not None and _pct(ti) is not None:
            k = _pct(ti)
            r = int(r + (255 - r) * (1 - k))
            g = int(g + (255 - g) * (1 - k))
            b = int(b + (255 - b) * (1 - k))
        r = max(0, min(255, r)); g = max(0, min(255, g)); b = max(0, min(255, b))
        return '{:02X}{:02X}{:02X}'.format(r, g, b)

    def _resolve_fill_hexes(sp_elem) -> list:
        """收集 shape (含子元素) 的 solidFill 颜色, 支持 srgbClr / schemeClr。返回 hex6 列表。"""
        out = []
        for sf in sp_elem.findall('.//a:solidFill', NS):
            srgb = sf.find('a:srgbClr', NS)
            schm = sf.find('a:schemeClr', NS)
            if srgb is not None and srgb.get('val'):
                out.append(_apply_lum_mod(srgb.get('val').upper(), srgb))
                continue
            if schm is not None and schm.get('val'):
                name: str = str(schm.get('val') or '')
                # bg1/tx1/bg2/tx2 是 clrMap 别名, 对应 lt1/dk1/lt2/dk2 (未解析 clrMap 时按常见默认)
                alias: str = {'bg1': 'lt1', 'tx1': 'dk1', 'bg2': 'lt2', 'tx2': 'dk2'}.get(name, name)
                base = theme_scheme.get(alias)
                if base:
                    out.append(_apply_lum_mod(base, schm))
        return out

    # ---------- (3) 圆形/椭圆几何判定 ----------
    def _is_circle_shape(shape) -> bool:
        sp = shape._element
        # 3a. prstGeom prst="ellipse" —— DrawingML 官方椭圆/圆形预设
        for pg in sp.findall('.//a:prstGeom', NS):
            if pg.get('prst') == 'ellipse':
                return True
        # 3b. python-pptx 的 auto_shape_type == OVAL (9)
        try:
            if getattr(shape, 'auto_shape_type', None) is not None and \
               shape.auto_shape_type == 9:
                return True
        except Exception:
            pass
        # 3c. 名称含 Oval/Ellipse/Circle
        name = (getattr(shape, 'name', '') or '').lower()
        if 'oval' in name or 'ellipse' in name or 'circle' in name:
            return True
        # 3d. custGeom + 近似 1:1 宽高比 (自由形状圆形): 兜底
        cust = sp.find('.//a:custGeom', NS)
        if cust is not None:
            try:
                w = shape.width; h = shape.height
                if w and h:
                    ratio = w / h if w >= h else h / w
                    if ratio <= 1.10:  # 宽高比 <= 1.10, 认作圆形
                        return True
            except Exception:
                pass
        return False

    # ---------- (4) 颜色分类: HSV + 显式容差 ----------
    #   容差策略 (与 rubric "橙/青/绿" 对齐):
    #     橙色: hue ∈ [15°, 45°], sat >= 0.35, val >= 0.55
    #     青色: hue ∈ [160°, 200°], sat >= 0.25, val >= 0.45
    #     绿色: hue ∈ [70°, 155°], sat >= 0.30, val >= 0.35
    def _hex_to_hsv(hex6: str):
        try:
            r = int(hex6[0:2], 16) / 255.0
            g = int(hex6[2:4], 16) / 255.0
            b = int(hex6[4:6], 16) / 255.0
        except Exception:
            return None
        mx = max(r, g, b); mn = min(r, g, b); d = mx - mn
        v = mx
        s = 0 if mx == 0 else d / mx
        if d == 0:
            h = 0.0
        elif mx == r:
            h = 60 * (((g - b) / d) % 6)
        elif mx == g:
            h = 60 * (((b - r) / d) + 2)
        else:
            h = 60 * (((r - g) / d) + 4)
        return (h, s, v)

    def _classify(hex6):
        hsv = _hex_to_hsv(hex6)
        if hsv is None:
            return None
        h, s, v = hsv
        if 15.0 <= h <= 45.0 and s >= 0.35 and v >= 0.55:
            return 'orange'
        if 160.0 <= h <= 200.0 and s >= 0.25 and v >= 0.45:
            return 'cyan'
        if 70.0 <= h <= 155.0 and s >= 0.30 and v >= 0.35:
            return 'green'
        return None

    # ---------- (5) 主循环 ----------
    out_of_bounds_pages = set()
    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not _is_circle_shape(shape):
                continue
            try:
                l = shape.left / 914400
                t = shape.top / 914400
                w = shape.width / 914400
                h = shape.height / 914400
            except Exception:
                continue
            hex_colors = _resolve_fill_hexes(shape._element)
            if not any(_classify(c) in ('orange', 'cyan', 'green') for c in hex_colors):
                continue
            # 超出边界: 任一边落在页面外 (给 0.01 in 容差, 避免浮点误差)
            eps = 0.01
            if l < -eps or t < -eps or (l + w) > sw_in + eps or (t + h) > sh_in + eps:
                out_of_bounds_pages.add(i)
    if out_of_bounds_pages:
        award(-5, "页面中橙/青/绿三色圆形图案超出边界",
              f"涉及页: {sorted(out_of_bounds_pages)}")
    else:
        log("  [ok] 三色圆形图案均未超出边界 (无扣分)")

    # -----------------------------------------------------------------------
    # 维度 2 评分完毕, 返回结构化命中项 (由 evaluate() 组装最终 dict)
    # -----------------------------------------------------------------------
    return True, "", hits


def evaluate(dir_path: str):
    """统一入口：接收"脚本所在目录的路径", 在该目录下定位并评估 PPT 文档。

    返回结构 (参见 "脚本接口差异与统一建议.md §2.2")::

        {
            "id": "072",
            "file_name": "xxx.pptx",
            "status": "ok" | "error",
            "error": None | str,
            "dim1_pass": bool,
            "dim1_reason": str,
            "dim2_items": [
                {"rule": <细则原文>, "max_delta": <满分>, "delta": <实际得分>,
                 "hit": <是否命中>, "detail": <可选说明>},
                ...
            ],
            "total_score": int,
            "max_score": int,
        }
    """
    # 满分 = 所有正分之和 (扣分项不计入满分)
    max_score = sum(d for d, _ in RULES if d > 0)
    base_result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": max_score,
    }

    try:
        pptx_path = _locate_pptx(dir_path)
        if pptx_path is None:
            base_result["status"] = "error"
            base_result["error"] = f"在目录 {dir_path} 下未找到 .pptx 文件"
            return base_result
        base_result["file_name"] = os.path.basename(pptx_path)

        dim1_pass, dim1_reason, hits = _evaluate_pptx(pptx_path)
        base_result["dim1_pass"] = dim1_pass
        base_result["dim1_reason"] = dim1_reason

        if not dim1_pass:
            # 维度一未通过 -> 维度二不评分, 总分 0
            base_result["dim2_items"] = []
            base_result["total_score"] = 0
            return base_result

        # 组装维度二逐项 (命中 + 未命中都返回)
        hit_map = {name: delta for delta, name, _detail in hits}
        dim2_items = []
        total_score = 0
        for max_delta, name in RULES:
            rule_text = RULE_TEXTS.get(name, name)
            if name in hit_map:
                actual_delta = hit_map[name]
                hit = True
            else:
                actual_delta = 0
                hit = False
            dim2_items.append({
                "rule": rule_text,
                "max_delta": max_delta,
                "delta": actual_delta,
                "hit": hit,
                "detail": "",
            })
            total_score += actual_delta

        base_result["dim2_items"] = dim2_items
        base_result["total_score"] = total_score
        return base_result

    except Exception as e:
        base_result["status"] = "error"
        base_result["error"] = f"{type(e).__name__}: {e}"
        return base_result


if __name__ == "__main__":
    # 本地调试: 默认以脚本所在目录作为 dir_path; 也可通过命令行覆盖
    _default_dir = os.path.dirname(os.path.abspath(__file__))
    _target_dir = sys.argv[1] if len(sys.argv) > 1 else _default_dir
    print(json.dumps(evaluate(_target_dir), ensure_ascii=False, indent=2))
