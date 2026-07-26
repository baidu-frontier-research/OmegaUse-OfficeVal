# -*- coding: utf-8 -*-
"""
PPT 自动评分脚本 officeval_060

对外统一接口：
    evaluate(dir_path: str) -> dict
        dir_path 为脚本所在目录，脚本自己在该目录中定位 .pptx 文档并评估。
        返回结构见项目《脚本接口差异与统一建议》§2.2。
"""
import io
import json
import os
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageFilter, ImageStat

SCRIPT_ID = "060"
EMU_PER_CM = 360000


def emu2cm(emu):
    return emu / EMU_PER_CM


def check_dim1(_prs, _results) -> bool:
    """维度1：可用与可修改性

    仅保留“文件可被 python-pptx 正常打开/解析”这一条底线——
    evaluate() 已在 Presentation() 打开失败时进入错误分支；到本函数时视为通过。
    其余关于页数=5、页面尺寸=16:9、第4页以外幻灯片正常、第4页对象无越界/裁切/重叠等
    门槛检查均按需求删除。参数保留仅为保持调用签名兼容。
    """
    return True


def get_slide4_shapes(slide, W_cm, H_cm):
    """收集第4页所有形状信息"""
    shapes = []
    for sh in slide.shapes:
        l = emu2cm(sh.left) if sh.left is not None else 0
        t = emu2cm(sh.top) if sh.top is not None else 0
        w = emu2cm(sh.width) if sh.width is not None else 0
        h = emu2cm(sh.height) if sh.height is not None else 0
        text = ""
        if sh.has_text_frame:
            text = " ".join(p.text for p in sh.text_frame.paragraphs).strip()
        shapes.append({"name": sh.name, "l": l, "t": t, "w": w, "h": h,
                       "r": l + w, "b": t + h, "text": text, "shape": sh})
    return shapes


def is_picture(sh):
    return sh["shape"].shape_type == MSO_SHAPE_TYPE.PICTURE


def load_shape_image(shape_info):
    """读取图片形状的像素；读取失败返回 None。"""
    try:
        blob = shape_info["shape"].image.blob
        return Image.open(io.BytesIO(blob)).convert("RGB")
    except Exception:
        return None


def column_content_bounds(img, white_threshold=245, gray_delta=8):
    """估算图片横向有内容的范围。

    将纯白/纯灰/近乎无纹理列视为空白列；返回 (left_px, right_px, content_ratio)。
    """
    if img is None:
        return None
    small = img.resize((400, max(1, int(400 * img.height / img.width))))
    pix = small.load()
    Wp, Hp = small.size
    has_content = []
    for x in range(Wp):
        content_pixels = 0
        for y in range(Hp):
            r, g, b = pix[x, y]
            bright = (r + g + b) / 3
            near_gray = max(r, g, b) - min(r, g, b) <= gray_delta
            # 非纯白/非纯灰，或亮度不高且有颜色差异 → 认为有视觉内容
            if not (bright >= white_threshold or (near_gray and bright >= 210)):
                content_pixels += 1
        has_content.append(content_pixels / Hp > 0.03)
    xs = [i for i, ok in enumerate(has_content) if ok]
    if not xs:
        return None
    left, right = min(xs), max(xs)
    return left, right, (right - left + 1) / Wp


def estimate_edge_blank_cm(img, slide_w_cm):
    bounds = column_content_bounds(img)
    if not bounds:
        return None
    left, right, coverage = bounds
    px_w = 400
    return left / px_w * slide_w_cm, (px_w - 1 - right) / px_w * slide_w_cm, coverage


# ── 维度二评分项定义（rule + max_delta），保持稳定顺序 ──────────────────────────
DIM2_RULES = [
    {"key": "bg_full_cover", "rule": "第4页整页背景图片：图片覆盖完整幻灯片，四边界距页面对应边缘0—0.3厘米", "max_delta": 5},
    {"key": "bg_subject_pos", "rule": "第4页背景图片主体位置：左侧高浊度水罐位于页面宽度0%—22%区域，中部膜堆位于30%—62%区域，且主要设备未被完全遮挡", "max_delta": 3},
    {"key": "bg_whitespace", "rule": "第4页背景图片留白控制：两侧无超过1厘米的纯白/纯灰/无内容区域，主体视觉内容覆盖页面宽90%以上", "max_delta": 3},
    {"key": "labels_missing", "rule": "第4页“预处理”“膜分离”“压差判读”“回用缓冲”4个标签缺少任意两个以上（扣分项）", "max_delta": -5},
]


def _new_dim2_items():
    """按 DIM2_RULES 生成默认未命中的 dim2_items 列表。"""
    items = []
    for r in DIM2_RULES:
        items.append({
            "rule": r["rule"],
            "max_delta": r["max_delta"],
            "delta": 0,
            "hit": False,
            "detail": "",
        })
    return items


def _dim2_item(items, key):
    """按 key 从 DIM2_RULES 找到 items 中对应条目。"""
    idx = next(i for i, r in enumerate(DIM2_RULES) if r["key"] == key)
    return items[idx]


def check_dim2(prs):
    """维度二：返回 dim2_items 列表（命中与未命中都在其中）。"""
    from PIL import ImageFilter, ImageStat  # 局部导入避免全局告警
    slide4 = prs.slides[3]
    W = emu2cm(prs.slide_width)   # 33.87
    H = emu2cm(prs.slide_height)  # 19.05
    shapes = get_slide4_shapes(slide4, W, H)
    items = _new_dim2_items()

    # ── 加分项 ──────────────────────────────────────────────────────────────

    # +5: 第4页整页背景图片：图片覆盖完整幻灯片，
    #     左边界位于距页面左边0—0.3厘米范围内，右边界位于距页面右边0—0.3厘米范围内，
    #     上边界位于距页面上边0—0.3厘米范围内，下边界位于距页面下边0—0.3厘米范围内。
    pics = [s for s in shapes if is_picture(s)]
    bg = None
    if pics:
        # 取覆盖面积最大的图片作为整页背景图片
        bg = max(pics, key=lambda s: s["w"] * s["h"])

    # 四边到页面对应边缘的距离（cm）
    #   左边界距页面左边 = 图片left
    #   右边界距页面右边 = W - 图片right
    #   上边界距页面上边 = 图片top
    #   下边界距页面下边 = H - 图片bottom
    it = _dim2_item(items, "bg_full_cover")
    if bg:
        left_dist = bg["l"]
        right_dist = W - bg["r"]
        top_dist = bg["t"]
        bottom_dist = H - bg["b"]
        # 每一边都必须落在 0—0.3 厘米范围内（含边界）
        if (0 <= left_dist <= 0.3 and 0 <= right_dist <= 0.3
                and 0 <= top_dist <= 0.3 and 0 <= bottom_dist <= 0.3):
            it["hit"] = True
            it["delta"] = 5
            it["detail"] = "背景图片覆盖完整"
        else:
            it["detail"] = "背景图片未完整覆盖页面"
    else:
        it["detail"] = "无背景图片"

    bg_img = load_shape_image(bg) if bg else None

    # +3: 第4页背景图片主体位置：
    #     左侧高浊度水罐位于页面宽度0%—22%区域，
    #     中部膜堆位于页面宽度30%—62%区域，
    #     主要设备没有被页面边缘或信息卡片完全遮挡。
    #
    # 背景图片为整页铺满的单张照片，水罐/膜堆是图片内的像素主体而非独立形状；
    # 由于背景图片铺满整页，页面宽度百分比与图片横向像素百分比一一对应。
    # 通过边缘/纹理密度确认对应区域确实存在设备主体，
    # 并检查信息卡片是否把该区域在页面上“完全遮挡”。
    it = _dim2_item(items, "bg_subject_pos")
    if bg and bg_img:
        small = bg_img.resize((400, max(1, int(400 * bg_img.height / bg_img.width)))).convert("L")
        edges = small.filter(ImageFilter.FIND_EDGES)
        Wp, Hp = edges.size

        def has_object(x0, x1):
            """页面宽度[x0,x1]区域内是否存在设备主体（边缘密度足够）。"""
            crop = edges.crop((int(x0 * Wp), 0, int(x1 * Wp), Hp))
            return ImageStat.Stat(crop).mean[0] > 2.0

        # 点1：左侧高浊度水罐位于页面宽度0%—22%区域
        left_tank_ok = has_object(0.00, 0.22)
        # 点2：中部膜堆位于页面宽度30%—62%区域
        mid_stack_ok = has_object(0.30, 0.62)

        # 点3：主要设备（左侧水罐 0%—22%、中部膜堆 30%—62%）没有被信息卡片完全遮挡。
        # “完全遮挡”指某个信息卡片在垂直和水平方向都完整盖住整个主体区域。
        def fully_covered(zx0, zx1):
            zl, zr = W * zx0, W * zx1
            for s in shapes:
                if is_picture(s):
                    continue
                # 卡片水平方向盖住整个区域，且垂直方向从区域上方延伸到下方（视为完全遮挡）
                if (s["l"] <= zl and s["r"] >= zr
                        and s["t"] <= 0 and s["b"] >= H):
                    return True
            return False

        blocked = fully_covered(0.00, 0.22) or fully_covered(0.30, 0.62)

        if left_tank_ok and mid_stack_ok and not blocked:
            it["hit"] = True
            it["delta"] = 3
            it["detail"] = "背景图片主体位置合理"
        else:
            it["detail"] = "背景图片主体位置/遮挡不符"
    else:
        it["detail"] = "无背景图片或图片不可读"

    # +3: 第4页背景图片留白控制：
    #     页面左侧和右侧没有连续宽度超过1厘米的纯白、纯灰或无内容区域；
    #     页面主体视觉内容覆盖页面宽90%以上。
    #
    # 背景图片铺满整页，图片横向像素百分比与页面宽度一一对应，
    # 因此对背景图片逐列检测纯白/纯灰/无内容列，得到：
    #   左侧连续空白宽度 = 图片最左侧连续空白列宽度
    #   右侧连续空白宽度 = 图片最右侧连续空白列宽度
    #   主体视觉内容覆盖宽度 = (最右有内容列 - 最左有内容列) / 页面宽
    it = _dim2_item(items, "bg_whitespace")
    blanks = estimate_edge_blank_cm(bg_img, W) if bg_img else None
    if blanks:
        left_blank, right_blank, coverage = blanks
        # 点1：左侧和右侧连续空白宽度都不超过1厘米
        edge_blank_ok = left_blank <= 1.0 and right_blank <= 1.0
        # 点2：主体视觉内容覆盖页面宽90%以上
        coverage_ok = coverage >= 0.90
        if edge_blank_ok and coverage_ok:
            it["hit"] = True
            it["delta"] = 3
            it["detail"] = "背景图片留白控制良好"
        else:
            it["detail"] = "留白控制不足"
    else:
        it["detail"] = "无背景图片或图片不可读，无法检查留白"

    # ── 扣分项 ──────────────────────────────────────────────────────────────

    required_labels = ["预处理", "膜分离", "压差判读", "回用缓冲"]

    # -5: 第4页“预处理”“膜分离”“压差判读”“回用缓冲”4个标签缺少任意两个以上。
    # “缺少任意两个以上”= 4个标签中缺失≥2个时扣分。
    # 逐个在页面全部文本框文字中查找该标签是否出现。
    all_text = " ".join(s["text"] for s in shapes if s["text"])
    missing_labels = [lb for lb in required_labels if lb not in all_text]
    it = _dim2_item(items, "labels_missing")
    if len(missing_labels) >= 2:
        it["hit"] = True
        it["delta"] = -5
        it["detail"] = f"缺失标签：{','.join(missing_labels)}"

    # 保留 detail 字段结构，但对外输出统一置空（不影响评分与命中判定）
    for _it in items:
        _it["detail"] = ""

    return items


def _find_target_file(dir_path):
    """在给定目录里定位待评估的 .pptx 文档。"""
    import os
    if not os.path.isdir(dir_path):
        return None
    candidates = [f for f in os.listdir(dir_path)
                  if f.lower().endswith(".pptx") and not f.startswith("~$")]
    if not candidates:
        return None
    # 若存在多个 .pptx，优先返回名字与脚本编号相关或最新修改的一个
    candidates.sort(key=lambda f: os.path.getmtime(os.path.join(dir_path, f)), reverse=True)
    return os.path.join(dir_path, candidates[0])


def evaluate(dir_path: str) -> dict:
    """统一入口：接收脚本所在目录，自行定位并评估其中的 .pptx 文档。"""
    import os
    from pptx import Presentation

    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": _new_dim2_items(),
        "total_score": 0,
        "max_score": sum(r["max_delta"] for r in DIM2_RULES if r["max_delta"] > 0),
    }

    try:
        target = _find_target_file(dir_path)
        if not target:
            result["status"] = "error"
            result["error"] = f"目录 {dir_path} 下未找到 .pptx 文档"
            return result
        result["file_name"] = os.path.basename(target)

        try:
            prs = Presentation(target)
        except Exception as e:
            result["status"] = "error"
            result["error"] = f"文件无法打开: {e}"
            return result

        dim1_reasons = []
        dim1_ok = check_dim1(prs, dim1_reasons)
        result["dim1_pass"] = dim1_ok
        if not dim1_ok:
            result["dim1_reason"] = "；".join(dim1_reasons)
            # 维度一未通过：dim2_items 保留结构但不评估，总分为0
            return result

        items = check_dim2(prs)
        result["dim2_items"] = items
        result["total_score"] = sum(i["delta"] for i in items)
        return result
    except Exception as e:
        result["status"] = "error"
        result["error"] = f"脚本运行异常: {e}"
        return result


if __name__ == "__main__":
    import json
    import os
    import sys

    # 本地调试：默认使用脚本所在目录；也可用命令行参数覆盖。
    default_dir = os.path.dirname(os.path.abspath(__file__))
    target_dir = sys.argv[1] if len(sys.argv) > 1 else default_dir
    print(json.dumps(evaluate(target_dir), ensure_ascii=False, indent=2))
