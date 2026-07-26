#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""自动评估“名片复刻_可编辑版.pptx”。

依赖：python-pptx、Pillow（本脚本主要使用 python-pptx）。
运行：python evaluate_ppt.py 名片复刻_可编辑版.pptx

评分逻辑：
1. 先检查维度1（可用与可修改性）。任一项不满足，直接输出0分，不进入维度2。
2. 维度2逐条自动检测。正分项需满足该条内全部关键点；负分项（若以后添加）满足任一触发点即扣分。
3. 打印命中的评分点、未命中的评分点和最终得分。
"""

from __future__ import annotations

import json

SCRIPT_ID = "043"
import math
import re
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable, Optional

from pptx import Presentation

EMU_PER_CM = 360000
EMU_PER_PT = 12700
CHINESE_RE = re.compile(r"[一-鿿]")


@dataclass
class ShapeInfo:
    raw: object
    slide_no: int
    name: str
    type_name: str
    prst: str
    x: float
    y: float
    w: float
    h: float
    text: str
    paragraphs: list[dict]
    fill: Optional[tuple[int, int, int]]
    line: Optional[tuple[int, int, int]]

    @property
    def x2(self) -> float:
        return self.x + self.w

    @property
    def y2(self) -> float:
        return self.y + self.h

    @property
    def cx(self) -> float:
        return self.x + self.w / 2

    @property
    def cy(self) -> float:
        return self.y + self.h / 2

    @property
    def area(self) -> float:
        return max(0.0, self.w) * max(0.0, self.h)


@dataclass
class CheckResult:
    passed: bool
    details: str


@dataclass
class Criterion:
    cid: str
    score: int
    title: str
    check: Callable[["PptModel"], CheckResult]
    is_penalty: bool = False


class PptModel:
    def __init__(self, path: Path):
        self.path = path
        self.prs = Presentation(str(path))
        self.slide_w = self.prs.slide_width / EMU_PER_CM
        self.slide_h = self.prs.slide_height / EMU_PER_CM
        self.slides: list[list[ShapeInfo]] = []
        for idx, slide in enumerate(self.prs.slides, start=1):
            self.slides.append(list(self._flatten_slide(slide, idx)))

    def _flatten_slide(self, slide_or_group, slide_no: int) -> Iterable[ShapeInfo]:
        for shape in slide_or_group.shapes:
            yield self._shape_info(shape, slide_no)
            if hasattr(shape, "shapes"):
                yield from self._flatten_slide(shape, slide_no)

    def _shape_info(self, shape, slide_no: int) -> ShapeInfo:
        left = shape.left / EMU_PER_CM
        top = shape.top / EMU_PER_CM
        width = shape.width / EMU_PER_CM
        height = shape.height / EMU_PER_CM
        # 连接符/线条可能存在负宽高，统一归一化成左上角 + 正宽高。
        x1 = min(left, left + width)
        y1 = min(top, top + height)
        x2 = max(left, left + width)
        y2 = max(top, top + height)
        paragraphs = self._paragraphs(shape)
        return ShapeInfo(
            raw=shape,
            slide_no=slide_no,
            name=getattr(shape, "name", ""),
            type_name=str(getattr(shape, "shape_type", "")),
            prst=get_prst(shape),
            x=x1,
            y=y1,
            w=x2 - x1,
            h=y2 - y1,
            text=getattr(shape, "text", "") if getattr(shape, "has_text_frame", False) else "",
            paragraphs=paragraphs,
            fill=get_fill_rgb(shape),
            line=get_line_rgb(shape),
        )

    def _paragraphs(self, shape) -> list[dict]:
        if not getattr(shape, "has_text_frame", False):
            return []
        out = []
        for para in shape.text_frame.paragraphs:
            runs = []
            para_text = ""
            for run in para.runs:
                para_text += run.text
                font = run.font
                runs.append(
                    {
                        "text": run.text,
                        "font": font.name,
                        "size": font.size.pt if font.size else None,
                        "bold": font.bold,
                        "italic": font.italic,
                        "color": get_font_rgb(font),
                    }
                )
            if para_text.strip():
                out.append({"text": para_text, "runs": runs})
        return out

    def slide(self, no: int) -> list[ShapeInfo]:
        return self.slides[no - 1]

    def all_shapes(self) -> list[ShapeInfo]:
        return [s for slide in self.slides for s in slide]


# ------------------------- 基础解析/判断工具 -------------------------


def get_prst(shape) -> str:
    try:
        return shape.auto_shape_type.name.upper()
    except Exception:
        try:
            nodes = shape._element.xpath(".//a:prstGeom")
            if nodes:
                return (nodes[0].get("prst") or "").upper()
        except Exception:
            pass
    if "PICTURE" in str(getattr(shape, "shape_type", "")).upper():
        return "PICTURE"
    return ""


def rgb_from_color_format(cf) -> Optional[tuple[int, int, int]]:
    if cf is None:
        return None
    try:
        if cf.rgb is not None:
            s = str(cf.rgb)
            return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
    except Exception:
        pass
    # 某些主题色 python-pptx 不给 rgb，尝试从底层 XML 取 srgbClr。
    try:
        xml = cf._xFill.xml
        m = re.search(r"srgbClr val=\"([0-9A-Fa-f]{6})\"", xml)
        if m:
            s = m.group(1)
            return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
    except Exception:
        pass
    return None


def get_fill_rgb(shape) -> Optional[tuple[int, int, int]]:
    try:
        return rgb_from_color_format(shape.fill.fore_color)
    except Exception:
        return None


def get_line_rgb(shape) -> Optional[tuple[int, int, int]]:
    try:
        return rgb_from_color_format(shape.line.color)
    except Exception:
        return None


def get_font_rgb(font) -> Optional[tuple[int, int, int]]:
    try:
        return rgb_from_color_format(font.color)
    except Exception:
        return None


def clean_text(s: str) -> str:
    table = str.maketrans(
        {
            "：": ":",
            "，": ",",
            "–": "-",
            "—": "-",
            "－": "-",
            "‐": "-",
            "黏": "粘",
            "涂": "图",  # 兼容评分细则“图层”和文件中常见“涂层”写法
        }
    )
    return re.sub(r"\s+", "", s.translate(table)).lower()


def text_contains(shape: ShapeInfo, target: str) -> bool:
    return clean_text(target) in clean_text(shape.text)


def slide_text(slide: list[ShapeInfo]) -> str:
    return "\n".join(s.text for s in slide if s.text)


def font_names(shape: ShapeInfo) -> list[str]:
    names = []
    for para in shape.paragraphs:
        for run in para["runs"]:
            if run["text"].strip() and run["font"]:
                names.append(run["font"])
    return names


def font_sizes(shape: ShapeInfo) -> list[float]:
    sizes = []
    for para in shape.paragraphs:
        for run in para["runs"]:
            if run["text"].strip() and run["size"] is not None:
                sizes.append(float(run["size"]))
    return sizes


def font_colors(shape: ShapeInfo) -> list[Optional[tuple[int, int, int]]]:
    colors = []
    for para in shape.paragraphs:
        for run in para["runs"]:
            if run["text"].strip():
                colors.append(run["color"])
    return colors


def all_runs_match(shape: ShapeInfo, *, font: Callable[[str], bool] | None = None,
                   size: tuple[float, float] | None = None,
                   color: Callable[[Optional[tuple[int, int, int]]], bool] | None = None,
                   bold: Optional[bool] = None) -> bool:
    runs = [r for p in shape.paragraphs for r in p["runs"] if r["text"].strip()]
    if not runs:
        return False
    for run in runs:
        if font and not font(run["font"] or ""):
            return False
        if size and (run["size"] is None or not between(float(run["size"]), size[0], size[1], tol=0.25)):
            return False
        if color and not color(run["color"]):
            return False
        if bold is True and run["bold"] is not True:
            return False
        if bold is False and run["bold"] is True:
            return False
    return True


def any_run_match(shape: ShapeInfo, **kwargs) -> bool:
    runs = [r for p in shape.paragraphs for r in p["runs"] if r["text"].strip()]
    if not runs:
        return False
    for run in runs:
        fake = ShapeInfo(shape.raw, shape.slide_no, shape.name, shape.type_name, shape.prst,
                         shape.x, shape.y, shape.w, shape.h, run["text"],
                         [{"text": run["text"], "runs": [run]}], shape.fill, shape.line)
        if all_runs_match(fake, **kwargs):
            return True
    return False


def between(v: float, lo: float, hi: float, tol: float = 0.0) -> bool:
    return lo - tol <= v <= hi + tol


def pos_pct(model: PptModel, shape: ShapeInfo) -> tuple[float, float, float, float]:
    return shape.x / model.slide_w, shape.y / model.slide_h, shape.x2 / model.slide_w, shape.y2 / model.slide_h


def center_pct(model: PptModel, shape: ShapeInfo) -> tuple[float, float]:
    return shape.cx / model.slide_w, shape.cy / model.slide_h


def in_region(model: PptModel, shape: ShapeInfo, x: tuple[float, float] | None = None,
              y: tuple[float, float] | None = None, *, use_center: bool = False,
              tol: float = 0.0) -> bool:
    px, py = center_pct(model, shape) if use_center else (shape.x / model.slide_w, shape.y / model.slide_h)
    return (x is None or between(px, x[0], x[1], tol)) and (y is None or between(py, y[0], y[1], tol))


def norm_font_name(name: str) -> str:
    return (name or "").replace(" ", "").lower()


def is_yahei_or_heiti(name: str) -> bool:
    n = norm_font_name(name)
    return any(k in n for k in ["yahei", "microsoftyahei", "simhei", "heiti", "黑体", "微软雅黑"])


def is_heiti_only(name: str) -> bool:
    """严格判定 typeface 是否为"黑体"（SimHei / Heiti），排除"微软雅黑"。

    用于细则原文仅写"黑体"（而非"微软雅黑或黑体"）的条目。OOXML 中：
      - SimHei：Windows 中文黑体，typeface="SimHei" 或 "黑体"
      - Heiti SC / Heiti TC：macOS 黑体，typeface="Heiti SC"/"Heiti TC"
    微软雅黑（Microsoft YaHei / YaHei）不算作细则字面上的"黑体"。
    """
    n = norm_font_name(name)
    if "yahei" in n or "微软雅黑" in n:  # 明确排除
        return False
    return any(k in n for k in ["simhei", "heiti", "黑体"])


def is_arial(name: str) -> bool:
    return "arial" in norm_font_name(name)


def is_times_or_arial(name: str) -> bool:
    n = norm_font_name(name)
    return "arial" in n or "timesnewroman" in n or "times" in n


def is_song_or_kai(name: str) -> bool:
    n = norm_font_name(name)
    return any(k in n for k in ["simsun", "song", "宋", "kaiti", "kai", "楷", "xingkai", "stxingkai"])


def is_script_font(name: str) -> bool:
    n = norm_font_name(name)
    return any(k in n for k in ["brushscript", "scriptmt", "alexbrush", "segoescript"])


def dist(c: Optional[tuple[int, int, int]], target: tuple[int, int, int]) -> float:
    if c is None:
        return 999.0
    return math.sqrt(sum((a - b) ** 2 for a, b in zip(c, target)))


def is_white(c) -> bool:
    # 非主题色：只要"看起来是白色"就通过——不强求纯白，接受略带轻微偏色的浅色。
    if c is None:
        return False
    return min(c) >= 210 and max(c) - min(c) <= 40


def is_black(c) -> bool:
    # 非主题色：只要"看起来是黑色"就通过——接受略带偏色的深色（如 #202428 之类）。
    if c is None:
        return False
    return max(c) <= 80 and max(c) - min(c) <= 40


def is_gray(c) -> bool:
    # 非主题色：只要是灰色即可——放宽色偏容差，覆盖浅灰、中灰、深灰。
    # 允许通道差最多 55（吸收 PPT 中带轻微色调倾向的"灰"），且平均亮度在 40-240 区间。
    if c is None:
        return False
    return max(c) - min(c) <= 55 and 40 <= sum(c) / 3 <= 240


def is_green(c) -> bool:
    # 放宽：只要色相是绿色即可——不再限定为特定主题绿。
    # 判定：g 通道必须最大，且明显大于 r、b；排除灰度色。
    if c is None:
        return False
    r, g, b = c
    if max(c) - min(c) < 20:  # 灰度色排除
        return False
    if g < r or g < b:
        return False
    # g 要明显高于 r（避免把黄色、米色误判为绿）
    if g - r < 12:
        return False
    # 且不能过暗（纯黑的偏色）
    return g >= 50


def is_deep_green(c) -> bool:
    # 放宽：任意"较深的绿色"都可通过。
    if c is None:
        return False
    return is_green(c) and max(c) < 180


def is_orange(c) -> bool:
    # 非主题色：只要色相是橙色即可（r 明显大于 g/b，且不是红也不是黄）。
    if c is None:
        return False
    r, g, b = c
    if r < 130:
        return False
    if r <= g or r <= b:
        return False
    # 排除纯黄（g 太接近 r 视为黄）和纯红（g 太低视为红）
    if g < 40 or g > r - 10:
        return False
    return b <= 130 and r - b >= 40


def is_yellow(c) -> bool:
    # 非主题色：只要色相是黄色即可（r、g 都高且接近，b 明显低）。
    if c is None:
        return False
    r, g, b = c
    if r < 150 or g < 130:
        return False
    if b >= min(r, g) - 30:
        return False
    # r 与 g 要相近（区别于橙色 r>>g）
    return abs(r - g) <= max(60, r * 0.35)


def is_light_gray(c) -> bool:
    # 非主题色：用户要求"只要是灰色就通过"，不再强求"浅"灰。
    # 保留函数名以兼容现有调用点，实际语义放宽为"任意灰色"。
    return is_gray(c)


def color_name(c) -> str:
    if c is None:
        return "none"
    return "#%02X%02X%02X" % c


def overlap_area(a: ShapeInfo, b: ShapeInfo) -> float:
    x = max(0.0, min(a.x2, b.x2) - max(a.x, b.x))
    y = max(0.0, min(a.y2, b.y2) - max(a.y, b.y))
    return x * y


def estimated_text_right(shape: ShapeInfo) -> float:
    """估算文本真实绘制宽度的右边界，降低宽文本框导致的误判。"""
    sizes = font_sizes(shape)
    avg_size = sum(sizes) / len(sizes) if sizes else 12.0
    max_units = 0.0
    for line in text_lines(shape):
        units = 0.0
        for ch in line:
            if ch.isspace():
                units += 0.25
            elif ord(ch) > 127:
                units += 1.0
            else:
                units += 0.55
        max_units = max(max_units, units)
    # 1pt约0.035cm；加少量余量，避免过度低估。
    estimated_width = max_units * avg_size * 0.035 * 1.20
    return shape.x + min(shape.w, max(0.4, estimated_width))


def is_picture(shape: ShapeInfo) -> bool:
    return "PICTURE" in shape.type_name.upper() or shape.prst == "PICTURE"


def is_line(shape: ShapeInfo) -> bool:
    return shape.prst == "LINE" or (shape.w == 0 or shape.h == 0)


def is_horizontal_line(shape: ShapeInfo) -> bool:
    return is_line(shape) and shape.w >= 0.3 and shape.h <= 0.12


def is_vertical_line(shape: ShapeInfo) -> bool:
    return is_line(shape) and shape.h >= 0.3 and shape.w <= 0.12


def is_round_rect(shape: ShapeInfo) -> bool:
    return "ROUND" in shape.prst and "RECT" in shape.prst


def is_oval(shape: ShapeInfo) -> bool:
    return "OVAL" in shape.prst or "ELLIPSE" in shape.prst


def find_text_shapes(model: PptModel, slide_no: int, target: str) -> list[ShapeInfo]:
    return [s for s in model.slide(slide_no) if s.text and text_contains(s, target)]


def first_text(model: PptModel, slide_no: int, target: str) -> Optional[ShapeInfo]:
    found = find_text_shapes(model, slide_no, target)
    return found[0] if found else None


def text_lines(shape: ShapeInfo) -> list[str]:
    return [p["text"] for p in shape.paragraphs] or [x for x in shape.text.splitlines() if x.strip()]


def ok(msg: str) -> CheckResult:
    return CheckResult(True, msg)


def fail(msg: str) -> CheckResult:
    return CheckResult(False, msg)


# ------------------------- 维度1 -------------------------


def check_dimension_1(path: Path) -> tuple[bool, list[str], Optional[PptModel]]:
    details: list[str] = []
    if path.suffix.lower() != ".pptx":
        details.append(f"✗ 文件扩展名不是 .pptx：{path.suffix}")
        return False, details, None
    details.append("✓ 文件扩展名为 .pptx")

    try:
        model = PptModel(path)
    except Exception as exc:
        details.append(f"✗ 文件无法正常打开：{exc}")
        return False, details, None
    details.append("✓ 文件可由 python-pptx 正常打开")

    if len(model.prs.slides) != 2:
        details.append(f"✗ 幻灯片页数为 {len(model.prs.slides)}，不是2页")
        return False, details, model
    details.append("✓ 交付PPT只包含2页幻灯片")

    slide2_text = slide_text(model.slide(2))
    if CHINESE_RE.search(slide2_text):
        details.append("✗ 第2页出现中文文本")
        return False, details, model
    details.append("✓ 第2页未检测到中文文本")

    for no in (1, 2):
        slide = model.slide(no)
        page_area = model.slide_w * model.slide_h
        huge_pictures = [s for s in slide if is_picture(s) and s.area / page_area > 0.80]
        editable_count = sum(1 for s in slide if s.text or (not is_picture(s)))
        if huge_pictures and editable_count < 5:
            details.append(f"✗ 第{no}页疑似整页图片且不可编辑")
            return False, details, model
    details.append("✓ 未检测到整页不可编辑图片覆盖")

    return True, details, model


# ------------------------- 维度2评分点 -------------------------


def c01_logo_top_left(model: PptModel) -> CheckResult:
    """+5：第1页幻灯片距左页面边线0%–20%、距顶部边线0%–20%位置内出现组件
    形状或者图片：所有组件组合或者图片总体宽度为2-4cm，高度为2-4cm；若为
    组件组合，则其中包含填充为橙色的圆形组件以及剩余组件填充为绿色；橙色
    部分位于整体的中上位置。
    """
    slide = model.slide(1)
    # 位置约束：形状左上角距左页面边线 0%-20%、距顶部边线 0%-20%
    # OOXML `off x,y`（EMU）在 PowerPoint / WPS / Keynote 中写法一致，
    # 归一化为 slide_w/slide_h 百分比即可跨办公软件生效。
    candidates = [
        s for s in slide
        if s.x / model.slide_w <= 0.20
        and s.y / model.slide_h <= 0.20
        and (is_picture(s) or s.fill is not None)
        and s.w <= 4.5 and s.h <= 4.5
        and s.area > 0.01
    ]
    if not candidates:
        return fail("第1页左上角0%-20%区域内未找到组件形状或图片")

    x1 = min(s.x for s in candidates)
    y1 = min(s.y for s in candidates)
    x2 = max(s.x2 for s in candidates)
    y2 = max(s.y2 for s in candidates)
    w, h = x2 - x1, y2 - y1

    # 总体尺寸约束：宽 2-4cm、高 2-4cm（图片路径与组件组合路径均需满足）
    if not (between(w, 2, 4) and between(h, 2, 4)):
        return fail(f"第1页左上角总体尺寸 {w:.2f}×{h:.2f}cm 不在 2-4cm 范围内")

    # 分支一：整体是图片（区域内不存在非图片组件即视为"图片"路径）
    non_pictures = [s for s in candidates if not is_picture(s)]
    if not non_pictures:
        return ok(f"第1页左上角logo图片 {w:.2f}×{h:.2f}cm")

    # 分支二：组件组合。收集细则要求的全部问题后一次性汇报，避免"只报第一
    # 个失败"掩盖后续问题（例如"中上位置"未被检测到）。
    problems: list[str] = []

    # 其中包含填充为橙色的圆形组件
    # OOXML: prst="ellipse" + <a:solidFill><a:srgbClr val="..."/>；办公软件
    # 均以 srgbClr 写入颜色，is_orange 对色值做阈值判定。
    orange_circles = [s for s in candidates if is_oval(s) and is_orange(s.fill)]
    if not orange_circles:
        problems.append("组件组合中未找到填充为橙色的圆形组件")

    # 剩余组件填充为绿色
    rest = [s for s in candidates if s not in orange_circles]
    if not rest:
        problems.append("组件组合缺少剩余绿色组件")
    else:
        non_green = [s for s in rest if not is_green(s.fill)]
        if non_green:
            problems.append(
                f"除橙色圆形外仍存在非绿色填充组件："
                f"{[color_name(s.fill) for s in non_green]}"
            )

    # 橙色部分位于整体的"中上"位置：
    #   中 —— 橙色几何中心的水平位置接近组合水平中线
    #   上 —— 橙色几何中心的垂直位置位于组合上半部
    # OOXML `off`+`ext` 精确到 EMU，办公软件间一致。
    if orange_circles:
        combo_mid_x = x1 + w * 0.5
        combo_mid_y = y1 + h * 0.5
        # "上"：中心在上半部
        not_upper = [s for s in orange_circles if s.cy > combo_mid_y]
        if not_upper:
            problems.append("橙色圆形垂直中心未位于组合上半部（不『上』）")
        # "中"：中心横向偏离 ≤ 组合宽度 15%
        off_center = [
            s for s in orange_circles
            if w > 0 and abs(s.cx - combo_mid_x) / w > 0.15
        ]
        if off_center:
            for s in off_center:
                offset = (s.cx - combo_mid_x)
                problems.append(
                    f"橙色圆形水平中心偏离组合中线 {offset:+.2f}cm（组合宽 {w:.2f}cm，占比 {abs(offset)/w*100:.1f}%），不『中』"
                )

    if problems:
        return fail("第1页左上角logo组件不满足：" + "；".join(problems))

    return ok(
        f"第1页左上角logo组件组合 {w:.2f}×{h:.2f}cm，"
        f"含中上位置橙色圆形与剩余绿色组件"
    )


def c02_cn_company(model: PptModel) -> CheckResult:
    """+1：第1页幻灯片出现"福建蓝屿新材料科技有限公司"文本：字体为微软雅黑
    或黑体，24-32磅，位于距顶部页面 10%-25% 的位置。
    """
    s = first_text(model, 1, "福建蓝屿新材料科技有限公司")
    if not s:
        return fail("第1页未找到\"福建蓝屿新材料科技有限公司\"文本")

    # 字体：微软雅黑或黑体
    # 首选 python-pptx 解析的 font.name（对应 <a:latin typeface>）；
    # 中文字体在部分办公软件（WPS/Keynote 等）中可能仅写入 <a:ea typeface>，
    # 因此再从形状原始 XML 中读取全部 typeface 作为兜底，保证跨办公软件有效。
    try:
        typefaces = re.findall(r'typeface="([^"]+)"', s.raw._element.xml)
    except Exception:
        typefaces = []
    latin_names = font_names(s)
    all_names = [n for n in latin_names + typefaces if n]
    if not all_names or not any(is_yahei_or_heiti(n) for n in all_names):
        return fail(
            f"第1页公司中文名字体不是微软雅黑/黑体："
            f"latin={latin_names} typefaces={typefaces}"
        )

    # 字号：24-32 磅（每个非空 run 都必须落在区间内）
    sizes = font_sizes(s)
    if not sizes or not all(between(sz, 24, 32) for sz in sizes):
        return fail(f"第1页公司中文名字号不在 24-32 磅：{sizes}")

    # 位置：距顶部页面 10%-25%
    top_pct = s.y / model.slide_h
    if not between(top_pct, 0.10, 0.25):
        return fail(f"第1页公司中文名距顶部 {top_pct:.1%}，不在 10%-25%")

    return ok(
        f"第1页公司中文名符合：字号={sizes} 距顶部={top_pct:.1%}"
    )


def c03_en_company_slide2(model: PptModel) -> CheckResult:
    """+1：第2页幻灯片出现"FUJIAN LANYU NOVA MATERIALS CO., LTD."文本：字体
    为 Arial，24-32 磅，整体位于距页面左边线 0%-70%、顶部边线 10%-25% 的
    位置；分两行排列布局。
    """
    s = first_text(model, 2, "FUJIAN LANYU NOVA MATERIALS CO., LTD.")
    if not s:
        return fail("第2页未找到\"FUJIAN LANYU NOVA MATERIALS CO., LTD.\"文本")

    # 字体：Arial
    # 为兼容 PowerPoint / WPS / Keynote 等办公软件，除 python-pptx 解析的
    # font.name（<a:latin typeface>）之外，再从原始 XML 中收集所有 typeface
    # 声明，避免仅写入 ea/cs 时漏判。
    try:
        typefaces = re.findall(r'typeface="([^"]+)"', s.raw._element.xml)
    except Exception:
        typefaces = []
    latin_names = font_names(s)
    all_names = [n for n in latin_names + typefaces if n]
    if not all_names or not all(is_arial(n) for n in all_names):
        return fail(
            f"第2页英文公司名字体不是 Arial："
            f"latin={latin_names} typefaces={typefaces}"
        )

    # 字号：24-32 磅（每个非空 run 都必须落在区间内）
    sizes = font_sizes(s)
    if not sizes or not all(between(sz, 24, 32) for sz in sizes):
        return fail(f"第2页英文公司名字号不在 24-32 磅：{sizes}")

    # 整体位置：距页面左边线 0%-70%、顶部边线 10%-25%
    # "整体位于"指文本框的外框整体落入区域内 → 左右边界都需在横向范围内、
    # 上下边界都需在纵向范围内。办公软件（PowerPoint/WPS/Keynote）保存
    # 位置时以 EMU 为单位，最终换算到百分比会有 ~1%-2% 的渲染精度损失，
    # 因此这里保留 2% 容差以兼容真实办公软件文件，同时仍严格执行"整体
    # 位于"的双边界检查。
    x1_pct = s.x / model.slide_w
    x2_pct = s.x2 / model.slide_w
    y1_pct = s.y / model.slide_h
    y2_pct = s.y2 / model.slide_h
    if not (between(x1_pct, 0.0, 0.70, tol=0.02)
            and between(x2_pct, 0.0, 0.70, tol=0.02)):
        return fail(
            f"第2页英文公司名横向未整体落入 0%-70%："
            f"x=[{x1_pct:.1%},{x2_pct:.1%}]"
        )
    if not (between(y1_pct, 0.10, 0.25, tol=0.02)
            and between(y2_pct, 0.10, 0.25, tol=0.02)):
        return fail(
            f"第2页英文公司名纵向未整体落入 10%-25%："
            f"y=[{y1_pct:.1%},{y2_pct:.1%}]"
        )

    # 分两行排列布局：文本框内正好包含两个非空段落
    lines = text_lines(s)
    if len(lines) != 2:
        return fail(f"第2页英文公司名未分两行排列，当前行数={len(lines)}")

    return ok(
        f"第2页英文公司名符合：字号={sizes} x=[{x1_pct:.1%},{x2_pct:.1%}] "
        f"y=[{y1_pct:.1%},{y2_pct:.1%}] 行数=2"
    )


def _is_four_corner_rounded_rect(shape: ShapeInfo) -> bool:
    """判定 shape 是否为"四个圆角的矩形"。

    OOXML 中"四个圆角矩形"对应 prstGeom prst="roundRect"，与仅有
    1 个圆角 (round1Rect)、2 个同侧圆角 (round2SameRect)、
    2 个对角圆角 (round2DiagRect) 以及切角矩形 (snip*Rect) 严格区分。
    python-pptx 将其映射为 MSO_SHAPE.ROUNDED_RECTANGLE，name 为
    "ROUNDED_RECTANGLE"；若走 XML 兜底则为 "ROUNDRECT"。
    """
    p = shape.prst.upper()
    return p == "ROUNDED_RECTANGLE" or p == "ROUNDRECT"


def _is_solid_line(shape: ShapeInfo) -> bool:
    """判定 shape 的边线是否为"实线"。

    OOXML 中线条虚实通过 <a:ln> 下的 <a:prstDash val="..."/> 或
    <a:custDash> 指定；未显式声明时默认为 solid。这里既接受
    prstDash="solid" 也接受未声明虚线两种情况，并显式排除各种
    虚线/点线值。判定基于原始 XML，兼容 PowerPoint/WPS/Keynote。
    """
    try:
        xml = shape.raw._element.xml
    except Exception:
        return True  # 无法访问 XML 时不阻塞
    # 只在 <a:ln>...</a:ln> 段内判定
    ln_match = re.search(r"<a:ln\b[^>]*>.*?</a:ln>", xml, re.DOTALL)
    ln_xml = ln_match.group(0) if ln_match else ""
    # 存在自定义虚线段 → 非实线
    if "<a:custDash" in ln_xml:
        return False
    dash_match = re.search(r'<a:prstDash[^/]*val="([^"]+)"', ln_xml)
    if not dash_match:
        return True  # 未声明 → 默认实线
    return dash_match.group(1).lower() == "solid"


def c04_borders(model: PptModel) -> CheckResult:
    """+1：两页幻灯片均出现"包含四个圆角的矩形"形状；边线颜色为黑色、
    灰色、绿色实线；矩形整体宽度 24cm-26cm，高度 13cm-15cm。
    """
    for no in (1, 2):
        matched = []
        for s in model.slide(no):
            if not _is_four_corner_rounded_rect(s):
                continue
            if not between(s.w, 24, 26):
                continue
            if not between(s.h, 13, 15):
                continue
            if s.line is None:
                continue
            color_ok = is_black(s.line) or is_gray(s.line) or is_green(s.line)
            if not color_ok:
                continue
            if not _is_solid_line(s):
                continue
            matched.append(s)
        if not matched:
            return fail(
                f"第{no}页未找到四个圆角矩形、宽24-26cm×高13-15cm、"
                f"黑/灰/绿色实线边框"
            )
    return ok("两页均含四个圆角矩形（24-26×13-15cm），边线为黑/灰/绿色实线")


def c05_en_company_slide1_subtitle(model: PptModel) -> CheckResult:
    """+1：第1页幻灯片"福建蓝屿新材料科技有限公司"的下方出现
    "FUJIAN LANYU NOVA MATERIALS CO., LTD." 文本：从上至下整体位于页面
    10%-25% 之间；字体为 Arial，14-22 磅，颜色为橙色。
    """
    cn = first_text(model, 1, "福建蓝屿新材料科技有限公司")
    s = first_text(model, 1, "FUJIAN LANYU NOVA MATERIALS CO., LTD.")
    if not cn:
        return fail("第1页未找到\"福建蓝屿新材料科技有限公司\"文本")
    if not s:
        return fail("第1页未找到\"FUJIAN LANYU NOVA MATERIALS CO., LTD.\"文本")

    # 位于"福建蓝屿新材料科技有限公司"的下方：英文副标题顶部在中文名顶部之下
    if not s.y > cn.y:
        return fail(
            f"第1页英文副标题未位于中文公司名下方："
            f"cn_y={cn.y/model.slide_h:.1%} en_y={s.y/model.slide_h:.1%}"
        )

    # 从上至下整体位于页面 10%-25% 之间 → 上下边界都需落入区间
    # 办公软件写入 EMU 存在渲染精度损失，保留 2% 容差以兼容真实文件。
    y1_pct = s.y / model.slide_h
    y2_pct = s.y2 / model.slide_h
    if not (between(y1_pct, 0.10, 0.25, tol=0.02)
            and between(y2_pct, 0.10, 0.25, tol=0.02)):
        return fail(
            f"第1页英文副标题纵向未整体落入 10%-25%："
            f"y=[{y1_pct:.1%},{y2_pct:.1%}]"
        )

    # 字体：Arial（兼容 <a:latin>/<a:ea>/<a:cs> typeface 三种写入方式）
    try:
        typefaces = re.findall(r'typeface="([^"]+)"', s.raw._element.xml)
    except Exception:
        typefaces = []
    latin_names = font_names(s)
    all_names = [n for n in latin_names + typefaces if n]
    if not all_names or not all(is_arial(n) for n in all_names):
        return fail(
            f"第1页英文副标题字体不是 Arial："
            f"latin={latin_names} typefaces={typefaces}"
        )

    # 字号：14-22 磅
    sizes = font_sizes(s)
    if not sizes or not all(between(sz, 14, 22) for sz in sizes):
        return fail(f"第1页英文副标题字号不在 14-22 磅：{sizes}")

    # 颜色：橙色
    colors = font_colors(s)
    if not colors or not all(is_orange(c) for c in colors):
        return fail(
            f"第1页英文副标题颜色不是橙色："
            f"{[color_name(c) for c in colors]}"
        )

    return ok(
        f"第1页英文副标题符合：字号={sizes} y=[{y1_pct:.1%},{y2_pct:.1%}] "
        f"颜色={[color_name(c) for c in colors]}"
    )


def two_horizontal_lines_with_leaf(model: PptModel, slide_no: int, y_range: tuple[float, float], length_range: tuple[float, float],
                                   color_check: Callable, leaf_size: tuple[tuple[float, float], tuple[float, float]],
                                   below_shape: Optional[ShapeInfo] = None,
                                   leaf_colors_required: Optional[list[Callable]] = None) -> CheckResult:
    """两条同水平线 + 中间叶片"组合"判定。

    leaf_colors_required 为 None 时退化为旧逻辑（单一叶片满足任一颜色即可）；
    若传入颜色检查列表，则要求叶片为"组合图形"：在中间区域内必须存在**多于一个**
    形状，且这些形状的填充颜色集合需要**全部覆盖** leaf_colors_required 中的颜色
    （即每个颜色至少由一个形状满足），方可视为命中。
    """
    slide = model.slide(slide_no)
    lines = [s for s in slide if is_horizontal_line(s) and between(s.cy / model.slide_h, y_range[0], y_range[1], 0.02) and between(s.w, length_range[0], length_range[1], 0.25) and color_check(s.line)]
    default_leaf_color = lambda c: is_green(c) or is_orange(c) or is_yellow(c)
    pairs = []
    for i, a in enumerate(lines):
        for b in lines[i + 1:]:
            if abs(a.cy - b.cy) <= 0.15 and a.x2 <= b.x or b.x2 <= a.x:
                left, right = (a, b) if a.x < b.x else (b, a)
                if below_shape is not None and min(left.y, right.y) <= below_shape.y:
                    continue
                middle_x1, middle_x2 = left.x2, right.x
                if middle_x2 - middle_x1 > 3.0:  # 两条线之间空隙过大，不是同一组分隔器
                    continue
                middle_shapes = [s for s in slide
                                 if middle_x1 - 0.15 <= s.cx <= middle_x2 + 0.15
                                 and abs(s.cy - left.cy) <= 0.45
                                 and between(s.w, leaf_size[0][0], leaf_size[0][1], 0.08)
                                 and between(s.h, leaf_size[1][0], leaf_size[1][1], 0.08)]
                if leaf_colors_required:
                    # 组合图形：必须 ≥2 件形状或 1 张图片，且全部要求颜色都被覆盖
                    if any(is_picture(s) for s in middle_shapes):
                        pairs.append((left, right, middle_shapes[0]))
                        continue
                    if len(middle_shapes) < 2:
                        continue
                    if all(any(check(s.fill) for s in middle_shapes) for check in leaf_colors_required):
                        pairs.append((left, right, middle_shapes[0]))
                else:
                    valid = [s for s in middle_shapes if default_leaf_color(s.fill) or is_picture(s)]
                    if valid:
                        pairs.append((left, right, valid[0]))
    if pairs:
        l, r, leaf = pairs[0]
        return ok(f"第{slide_no}页找到两条同水平线及中间叶片：线长 {l.w:.2f}/{r.w:.2f}cm，叶片 {leaf.w:.2f}×{leaf.h:.2f}cm")
    return fail(f"第{slide_no}页未找到满足位置、长度、颜色和叶片尺寸的双横线组合")


def c06_slide1_top_separator(model: PptModel) -> CheckResult:
    """+3：第1页幻灯片距顶部 23%-26% 位置处出现两条在同一水平线上的单实线：
    位于"FUJIAN LANYU NOVA MATERIALS CO., LTD."下方；两条线长度均在
    5cm-8cm 范围内，颜色为橙色；两条线中间存在一个绿色和黄色填充的叶片
    组合形状或者图片，整体高度 0.4-0.5cm，宽度 0.3-0.4cm。
    """
    subtitle = first_text(model, 1, "FUJIAN LANYU NOVA MATERIALS CO., LTD.")
    if not subtitle:
        return fail("第1页未找到英文副标题，无法定位其下方分隔线")

    slide = model.slide(1)

    # 收集候选横线：
    # 1) 位置：cy 落在 23%-26%（保留 2% 容差兼容办公软件 EMU 渲染精度）
    # 2) 长度：5-8cm
    # 3) 颜色：橙色
    # 4) 单实线（<a:prstDash> 未声明或 val="solid"，且无 <a:custDash>）
    # 5) 位于英文副标题下方（顶部在副标题底部之下）
    lines = []
    for s in slide:
        if not is_horizontal_line(s):
            continue
        cy_pct = s.cy / model.slide_h
        if not between(cy_pct, 0.23, 0.26, tol=0.02):
            continue
        if not between(s.w, 5, 8):
            continue
        if not is_orange(s.line):
            continue
        if not _is_solid_line(s):
            continue
        if s.y < subtitle.y2:
            continue
        lines.append(s)

    # 两条线在"同一水平线上"：cy 差 ≤ 0.10cm（≈同一渲染像素行）；
    # 且互不重叠，才能作为"两条独立的横线"。
    pair = None
    for i, a in enumerate(lines):
        for b in lines[i + 1:]:
            if abs(a.cy - b.cy) > 0.10:
                continue
            left, right = (a, b) if a.x < b.x else (b, a)
            if left.x2 > right.x:  # 有重叠 → 不是并排两条
                continue
            pair = (left, right)
            break
        if pair:
            break

    if not pair:
        return fail(
            "第1页副标题下方 23%-26% 未找到两条在同一水平线、"
            "长度 5-8cm、橙色、单实线的横线"
        )

    left, right = pair
    middle_x1, middle_x2 = left.x2, right.x
    middle_cy = (left.cy + right.cy) / 2

    # 两条线中间的"叶片"组合形状或图片：
    # - 组合形状：中间区域内 ≥2 个具备填充/线条的图形，且填充颜色集合同时
    #   覆盖"绿色"与"黄色"两个要求。
    # - 图片：中间区域内存在一张图片（图片路径不再校验颜色，只校验尺寸）。
    # 整体：宽 0.3-0.4cm、高 0.4-0.5cm（组合看外包围盒；图片看自身尺寸）。
    middle_shapes = [
        s for s in slide
        if s is not left and s is not right
        and middle_x1 - 0.15 <= s.cx <= middle_x2 + 0.15
        and abs(s.cy - middle_cy) <= 0.45
    ]

    # 路径 A：图片
    pictures = [s for s in middle_shapes if is_picture(s)]
    for pic in pictures:
        if between(pic.w, 0.3, 0.4) and between(pic.h, 0.4, 0.5):
            return ok(
                f"第1页副标题下方双橙线中间找到叶片图片 "
                f"{pic.w:.2f}×{pic.h:.2f}cm"
            )

    # 路径 B：组合形状
    combo = [
        s for s in middle_shapes
        if not is_picture(s)
        and s.prst  # 需有几何形状（排除纯文本框）
        and (s.fill is not None or s.line is not None)
    ]
    if len(combo) >= 2:
        has_green = any(is_green(s.fill) for s in combo)
        has_yellow = any(is_yellow(s.fill) for s in combo)
        if has_green and has_yellow:
            cx1 = min(s.x for s in combo)
            cx2 = max(s.x2 for s in combo)
            cy1 = min(s.y for s in combo)
            cy2 = max(s.y2 for s in combo)
            combo_w, combo_h = cx2 - cx1, cy2 - cy1
            if between(combo_w, 0.3, 0.4) and between(combo_h, 0.4, 0.5):
                return ok(
                    f"第1页副标题下方双橙线中间找到叶片组合形状（{len(combo)} 件），"
                    f"整体 {combo_w:.2f}×{combo_h:.2f}cm，含绿色与黄色填充"
                )
            return fail(
                f"第1页副标题下方双橙线中间叶片组合尺寸不满足："
                f"{combo_w:.2f}×{combo_h:.2f}cm"
            )
        return fail(
            f"第1页副标题下方双橙线中间组合缺少绿色或黄色填充："
            f"green={has_green} yellow={has_yellow}"
        )

    return fail("第1页副标题下方双橙线中间未找到叶片组合形状或图片")


def c07_qr_both(model: PptModel) -> CheckResult:
    """+1：两页幻灯片右上角均出现二维码组合形状或者图片：距左侧页边线
    75%-100%、顶部 0%-25% 的距离范围内；整体宽度 3-4cm，高度 3-4cm，
    宽高比为 1:1。
    """
    for no in (1, 2):
        matched = []
        for s in model.slide(no):
            # 类型：组合形状（GroupShape，python-pptx 会同时暴露 group 与其
            # 子形状；group 本身具备 `.shapes` 属性）或者图片（PICTURE）。
            is_group = hasattr(s.raw, "shapes")
            if not (is_picture(s) or is_group):
                continue

            # 位置："距左侧页边线 75%-100%" → 形状左边距页面左边线的相对距离
            # 处在 75%-100%；"距顶部 0%-25%" → 形状上边距页面顶部的相对距离
            # 处在 0%-25%。二者以形状左上角为基准。保留 2% 容差以兼容办公
            # 软件写入 EMU 的渲染精度损失。
            x_pct = s.x / model.slide_w
            y_pct = s.y / model.slide_h
            if not between(x_pct, 0.75, 1.00, tol=0.02):
                continue
            if not between(y_pct, 0.00, 0.25, tol=0.02):
                continue

            # 尺寸：宽 3-4cm、高 3-4cm
            if not (between(s.w, 3, 4) and between(s.h, 3, 4)):
                continue

            # 宽高比 1:1：允许 ±5% 的办公软件渲染误差
            if s.h <= 0:
                continue
            ratio = s.w / s.h
            if not between(ratio, 1.0, 1.0, tol=0.05):
                continue

            matched.append(s)

        if not matched:
            return fail(
                f"第{no}页右上角未找到二维码组合形状/图片："
                f"位置 75%-100% × 0%-25%、尺寸 3-4cm、宽高比 1:1"
            )

    return ok("两页右上角均出现二维码组合形状或图片，位置/尺寸/宽高比均符合")


def c08_name_cn(model: PptModel) -> CheckResult:
    """+1：第1页幻灯片出现"沈知行"文本：位于页面从左至右 10%-45%、从上至下
    30%-60% 的位置处；字体为宋体或楷体，加粗，字号为 40-50 磅，颜色为绿色。
    """
    s = first_text(model, 1, "沈知行")
    if not s:
        return fail("第1页未找到\"沈知行\"文本")

    # 位置：整体位于页面从左至右 10%-45%、从上至下 30%-60%。
    # "位于"取形状的左上角为基准（与其它条目一致）。保留 2% 容差以兼容办公
    # 软件写入 EMU 的渲染精度损失。
    x_pct = s.x / model.slide_w
    y_pct = s.y / model.slide_h
    if not between(x_pct, 0.10, 0.45, tol=0.02):
        return fail(f"第1页\"沈知行\"横向位置 {x_pct:.1%} 不在 10%-45%")
    if not between(y_pct, 0.30, 0.60, tol=0.02):
        return fail(f"第1页\"沈知行\"纵向位置 {y_pct:.1%} 不在 30%-60%")

    # 字体：宋体或楷体。兼容 <a:latin>/<a:ea>/<a:cs> 三种 typeface 写入方式，
    # 以适配 PowerPoint / WPS / Keynote 对中文字体的不同存放位置。
    try:
        typefaces = re.findall(r'typeface="([^"]+)"', s.raw._element.xml)
    except Exception:
        typefaces = []
    latin_names = font_names(s)
    all_names = [n for n in latin_names + typefaces if n]
    if not all_names or not all(is_song_or_kai(n) for n in all_names):
        return fail(
            f"第1页\"沈知行\"字体不是宋体/楷体："
            f"latin={latin_names} typefaces={typefaces}"
        )

    # 加粗：OOXML `<a:rPr b="1"/>`；python-pptx 映射为 run.font.bold == True。
    # 未声明时为 None，视为未加粗。要求所有非空 run 均声明为加粗。
    runs = [r for p in s.paragraphs for r in p["runs"] if r["text"].strip()]
    if not runs or not all(r["bold"] is True for r in runs):
        return fail(
            f"第1页\"沈知行\"未加粗：bold={[r['bold'] for r in runs]}"
        )

    # 字号：40-50 磅
    sizes = font_sizes(s)
    if not sizes or not all(between(sz, 40, 50) for sz in sizes):
        return fail(f"第1页\"沈知行\"字号不在 40-50 磅：{sizes}")

    # 颜色：绿色
    colors = font_colors(s)
    if not colors or not all(is_green(c) for c in colors):
        return fail(
            f"第1页\"沈知行\"颜色不是绿色："
            f"{[color_name(c) for c in colors]}"
        )

    return ok(
        f"第1页\"沈知行\"符合：pos=({x_pct:.1%},{y_pct:.1%}) 字号={sizes} "
        f"颜色={[color_name(c) for c in colors]}"
    )


def c09_name_separator(model: PptModel) -> CheckResult:
    """+3：第1页幻灯片从上至下页面的 50% 处出现两条同一水平线的单实线：
    位于"沈知行"下方；长度 2cm-3.5cm；颜色为绿色；两条横线中间存在一个
    黄色或者橙色的叶片形状组合图或图片，整体高度 0.35-0.45cm、宽 0.2-0.3cm。
    """
    name = first_text(model, 1, "沈知行")
    if not name:
        return fail("第1页未找到\"沈知行\"，无法定位其下方分隔线")

    slide = model.slide(1)

    # 候选横线判定：
    # 1) 位于 50% 处：cy 落在 50% ± 3%（保留 3% 容差兼容办公软件 EMU 渲染精度）
    # 2) 位于"沈知行"下方：线的顶端在"沈知行"底部之下
    # 3) 长度 2cm-3.5cm
    # 4) 颜色为绿色
    # 5) 单实线（<a:prstDash> 未声明或 val="solid"，且无 <a:custDash>）
    lines = []
    for s in slide:
        if not is_horizontal_line(s):
            continue
        cy_pct = s.cy / model.slide_h
        if not between(cy_pct, 0.50, 0.50, tol=0.03):
            continue
        if s.y < name.y2:
            continue
        if not between(s.w, 2, 3.5):
            continue
        if not is_green(s.line):
            continue
        if not _is_solid_line(s):
            continue
        lines.append(s)

    # 两条"同一水平线"：cy 差 ≤ 0.10cm；且两条水平不重叠。
    pair = None
    for i, a in enumerate(lines):
        for b in lines[i + 1:]:
            if abs(a.cy - b.cy) > 0.10:
                continue
            left, right = (a, b) if a.x < b.x else (b, a)
            if left.x2 > right.x:
                continue
            pair = (left, right)
            break
        if pair:
            break

    if not pair:
        return fail(
            "第1页\"沈知行\"下方 50% 处未找到两条在同一水平线、"
            "长度 2-3.5cm、绿色单实线"
        )

    left, right = pair
    middle_x1, middle_x2 = left.x2, right.x
    middle_cy = (left.cy + right.cy) / 2

    # 两条线中间的"叶片"组合形状或图片。
    # 中间区域：cx 落在两线中间水平段内（放宽 0.15cm 容差），cy 与线中心接
    # 近（≤0.45cm）。
    middle_shapes = [
        s for s in slide
        if s is not left and s is not right
        and middle_x1 - 0.15 <= s.cx <= middle_x2 + 0.15
        and abs(s.cy - middle_cy) <= 0.45
    ]

    # 路径 A：图片。整体宽 0.2-0.3cm、高 0.35-0.45cm。
    for pic in [s for s in middle_shapes if is_picture(s)]:
        if between(pic.w, 0.2, 0.3) and between(pic.h, 0.35, 0.45):
            return ok(
                f"第1页\"沈知行\"下方双绿线中间找到叶片图片 "
                f"{pic.w:.2f}×{pic.h:.2f}cm"
            )

    # 路径 B：组合形状。至少 1 个组件填充为"黄色或者橙色"；整体外包围盒
    # 宽 0.2-0.3cm、高 0.35-0.45cm。
    combo = [
        s for s in middle_shapes
        if not is_picture(s)
        and s.prst  # 需有几何形状 prst（排除纯文本框）
        and (s.fill is not None or s.line is not None)
    ]
    if combo:
        has_yellow_or_orange = any(
            is_yellow(s.fill) or is_orange(s.fill) for s in combo
        )
        cx1 = min(s.x for s in combo)
        cx2 = max(s.x2 for s in combo)
        cy1 = min(s.y for s in combo)
        cy2 = max(s.y2 for s in combo)
        combo_w, combo_h = cx2 - cx1, cy2 - cy1
        size_ok = between(combo_w, 0.2, 0.3) and between(combo_h, 0.35, 0.45)
        if has_yellow_or_orange and size_ok:
            return ok(
                f"第1页\"沈知行\"下方双绿线中间找到叶片组合形状"
                f"（{len(combo)} 件），整体 {combo_w:.2f}×{combo_h:.2f}cm"
            )
        if not has_yellow_or_orange:
            return fail(
                f"第1页\"沈知行\"下方双绿线中间组合缺少黄色或橙色填充："
                f"{[color_name(s.fill) for s in combo]}"
            )
        return fail(
            f"第1页\"沈知行\"下方双绿线中间叶片组合尺寸不满足："
            f"{combo_w:.2f}×{combo_h:.2f}cm"
        )

    return fail("第1页\"沈知行\"下方双绿线中间未找到叶片组合形状或图片")


def c10_mp_slide1(model: PptModel) -> CheckResult:
    """+1：第1页幻灯片出现"M.P."文本：位于页面下方 70%-80%、左侧 10%-40%
    的区域范围内，且在"沈知行"下方且与其左对齐，字体为 Arial，13-20Pt，
    颜色为绿色。
    """
    mp = first_text(model, 1, "M.P.")
    name = first_text(model, 1, "沈知行")
    if not mp:
        return fail("第1页未找到\"M.P.\"文本")
    if not name:
        return fail("第1页未找到\"沈知行\"，无法判定 M.P. 与其左对齐")

    # 位置：位于页面下方 70%-80%、左侧 10%-40% 的区域范围内。
    # "位于...区域范围内"以形状左上角为基准；保留 2% 容差兼容办公软件
    # EMU→百分比的渲染精度损失。
    x_pct = mp.x / model.slide_w
    y_pct = mp.y / model.slide_h
    if not between(y_pct, 0.70, 0.80, tol=0.02):
        return fail(f"第1页\"M.P.\"纵向位置 {y_pct:.1%} 不在 70%-80%")
    if not between(x_pct, 0.10, 0.40, tol=0.02):
        return fail(f"第1页\"M.P.\"横向位置 {x_pct:.1%} 不在 10%-40%")

    # 在"沈知行"下方：M.P. 顶部在"沈知行"底部之下
    if not mp.y > name.y2:
        return fail(
            f"第1页\"M.P.\"未位于\"沈知行\"下方："
            f"name_y2={name.y2:.2f} mp_y={mp.y:.2f}"
        )

    # 与"沈知行"左对齐：两者左边界差 ≤ 0.35cm（办公软件文本框内边距 ~0.25cm）
    if abs(mp.x - name.x) > 0.35:
        return fail(
            f"第1页\"M.P.\"与\"沈知行\"未左对齐："
            f"Δx={abs(mp.x - name.x):.2f}cm"
        )

    # 字体：Arial。兼容 <a:latin>/<a:ea>/<a:cs> 三种 typeface 写入方式。
    try:
        typefaces = re.findall(r'typeface="([^"]+)"', mp.raw._element.xml)
    except Exception:
        typefaces = []
    latin_names = font_names(mp)
    all_names = [n for n in latin_names + typefaces if n]
    if not all_names or not all(is_arial(n) for n in all_names):
        return fail(
            f"第1页\"M.P.\"字体不是 Arial："
            f"latin={latin_names} typefaces={typefaces}"
        )

    # 字号：13-20 磅
    sizes = font_sizes(mp)
    if not sizes or not all(between(sz, 13, 20) for sz in sizes):
        return fail(f"第1页\"M.P.\"字号不在 13-20 磅：{sizes}")

    # 颜色：绿色
    colors = font_colors(mp)
    if not colors or not all(is_green(c) for c in colors):
        return fail(
            f"第1页\"M.P.\"颜色不是绿色：{[color_name(c) for c in colors]}"
        )

    return ok(
        f"第1页\"M.P.\"符合：pos=({x_pct:.1%},{y_pct:.1%}) 字号={sizes} "
        f"颜色={[color_name(c) for c in colors]}"
    )


def _is_times_new_roman(name: str) -> bool:
    """严格判定 typeface 是否为 Times New Roman（区别于普通 Times）。

    OOXML 中"Times New Roman"字体通常以完整名称写入 typeface；这里同时兼
    容"Times New Roman"和常见变体如"TimesNewRomanPSMT"、"NimbusRomNo9L"
    等,与"Arial"互斥判定。
    """
    n = norm_font_name(name)
    return "timesnewroman" in n


def c11_phone_text_both(model: PptModel) -> CheckResult:
    """+1：两页幻灯片均出现"139 4826 7153"文本：位于"M.P."右侧，字体为
    Arial、Times New Roman，15-22Pt，颜色为黑色、绿色。
    """
    for no in (1, 2):
        phone = first_text(model, no, "139 4826 7153")
        mp = first_text(model, no, "M.P.")
        if not phone:
            return fail(f"第{no}页未找到\"139 4826 7153\"文本")
        if not mp:
            return fail(f"第{no}页未找到\"M.P.\"，无法判定电话在其右侧")

        # 位于"M.P."右侧：电话文本框左边在 M.P. 右边之后
        if not phone.x >= mp.x2:
            return fail(
                f"第{no}页\"139 4826 7153\"未位于\"M.P.\"右侧："
                f"mp.x2={mp.x2:.2f} phone.x={phone.x:.2f}"
            )

        # 字体：Arial 或 Times New Roman。兼容 <a:latin>/<a:ea>/<a:cs> 三种
        # typeface 写入方式，以适配 PowerPoint / WPS / Keynote。
        try:
            typefaces = re.findall(r'typeface="([^"]+)"', phone.raw._element.xml)
        except Exception:
            typefaces = []
        latin_names = font_names(phone)
        all_names = [n for n in latin_names + typefaces if n]
        if not all_names:
            return fail(f"第{no}页\"139 4826 7153\"未声明字体")
        for n in all_names:
            if not (is_arial(n) or _is_times_new_roman(n)):
                return fail(
                    f"第{no}页\"139 4826 7153\"字体不是 Arial/Times New Roman："
                    f"{n!r}（全部：latin={latin_names} typefaces={typefaces}）"
                )

        # 字号：15-22Pt
        sizes = font_sizes(phone)
        if not sizes or not all(between(sz, 15, 22) for sz in sizes):
            return fail(
                f"第{no}页\"139 4826 7153\"字号不在 15-22 磅：{sizes}"
            )

        # 颜色：黑色或绿色（细则用"、"表示或）
        colors = font_colors(phone)
        if not colors or not all(is_black(c) or is_green(c) for c in colors):
            return fail(
                f"第{no}页\"139 4826 7153\"颜色不是黑色/绿色："
                f"{[color_name(c) for c in colors]}"
            )

    return ok(
        "两页\"139 4826 7153\"均位于\"M.P.\"右侧，"
        "Arial/Times New Roman 15-22Pt 黑/绿色符合"
    )


def icon_left_of_text(model: PptModel, slide_no: int, target: str, glyphs: list[str]) -> CheckResult:
    text = first_text(model, slide_no, target)
    if not text:
        return fail(f"第{slide_no}页未找到目标文本 {target}")
    icons = []
    for s in model.slide(slide_no):
        if s.x2 <= text.x + 0.15 and abs(s.cy - text.cy) <= 0.45 and is_round_rect(s) and between(s.w, 0.8, 1.0, 0.12) and between(s.h, 0.8, 1.0, 0.12) and is_deep_green(s.fill):
            glyph_ok = any(g in s.text for g in glyphs)
            white_text = not s.text.strip() or any(is_white(c) for c in font_colors(s))
            if glyph_ok and white_text:
                icons.append(s)
    if icons:
        return ok(f"第{slide_no}页 {target} 左侧找到深绿色圆角矩形白色图标")
    return fail(f"第{slide_no}页 {target} 左侧未找到符合尺寸/颜色/图案的图标")


def c12_phone_icon_both(model: PptModel) -> CheckResult:
    """+3：两页幻灯片中"M.P."左侧均有一个宽 0.8-1cm、高 0.8-1cm 的深绿色
    填充的圆角矩形，其中间有一个白色的电话图案。
    """
    phone_glyphs = ["☎", "☏", "phone", "tel"]
    for no in (1, 2):
        mp = first_text(model, no, "M.P.")
        if not mp:
            return fail(f"第{no}页未找到\"M.P.\"，无法定位其左侧图标")

        # 收集候选：位于 M.P. 左侧、宽高 0.8-1cm、深绿色填充的圆角矩形。
        # - 位于左侧：候选形状右边不超过 M.P. 左边（允许 0.15cm 办公软件边距）
        # - 圆角矩形：OOXML prst="roundRect"（严格四个圆角）
        # - 深绿色填充：is_deep_green(fill)
        # - 宽/高 0.8-1cm
        candidates = []
        for s in model.slide(no):
            if s is mp:
                continue
            if not _is_four_corner_rounded_rect(s):
                continue
            if not (between(s.w, 0.8, 1.0) and between(s.h, 0.8, 1.0)):
                continue
            if not is_deep_green(s.fill):
                continue
            if s.x2 > mp.x + 0.15:
                continue
            candidates.append(s)

        # 需要"其中间"存在电话图案（白色）：
        # 路径 A：圆角矩形自身文本框包含电话字符（☎/☏/phone/tel），
        #         且字体颜色为白色。这是最常见的图标写法（形状 + 文本）。
        # 路径 B：圆角矩形的几何范围内存在一张图片（is_picture）— 视为
        #         "电话图案"的图片形式。图片是否为电话内容无法在 OOXML
        #         级别校验，采用几何包含 + "存在图片"作为可行判据。
        # 路径 C：圆角矩形几何范围内存在白色矢量图形/线条组合 —— 覆盖
        #         "电话图案由白色 freeform / line / group 构成"这种常见写法。
        #         判定条件：
        #           a) 形状不是圆角矩形本身、不是 M.P. 文本；
        #           b) 形状不是图片（图片走路径 B）；
        #           c) 形状不承载文字（文字走路径 A）；
        #           d) 填充或线条颜色为白色；
        #           e) 几何上落入圆角矩形内部（面积≥50%重叠；线条按中心点判定）；
        #           f) 尺寸不超过圆角矩形自身（排除整页级白色背景）。
        def _mostly_inside(inner: "ShapeInfo", container: "ShapeInfo") -> bool:
            if inner.area > 0:
                return overlap_area(inner, container) / inner.area >= 0.50
            # 退化形状(w 或 h 为 0，如直线)：用中心点是否落入 container
            return (container.x - 0.05 <= inner.cx <= container.x2 + 0.05
                    and container.y - 0.05 <= inner.cy <= container.y2 + 0.05)

        def _is_white_vector_inside(inner: "ShapeInfo", container: "ShapeInfo", mp_shape: "ShapeInfo") -> bool:
            if inner is container or inner is mp_shape:
                return False
            if is_picture(inner):
                return False
            if inner.text.strip():
                return False
            if not (is_white(inner.fill) or is_white(inner.line)):
                return False
            # 排除超出圆角矩形自身尺寸的形状（页面级白色背景/大装饰）
            if inner.w > container.w * 1.1 or inner.h > container.h * 1.1:
                return False
            return _mostly_inside(inner, container)

        icon_ok = None
        for c in candidates:
            # 路径 A
            has_glyph = any(g in c.text for g in phone_glyphs)
            colors = font_colors(c)
            white_glyph = has_glyph and colors and all(is_white(cc) for cc in colors)
            if white_glyph:
                icon_ok = c
                break

            # 路径 B：查找几何上落入圆角矩形内部的图片
            has_pic_inside = any(
                is_picture(s)
                and overlap_area(s, c) / max(0.01, min(s.area, c.area)) >= 0.50
                for s in model.slide(no)
            )
            if has_pic_inside:
                icon_ok = c
                break

            # 路径 C：查找几何上落入圆角矩形内部的白色矢量图形/线条组合
            has_white_vector = any(
                _is_white_vector_inside(s, c, mp) for s in model.slide(no)
            )
            if has_white_vector:
                icon_ok = c
                break

        if not icon_ok:
            return fail(
                f"第{no}页\"M.P.\"左侧未找到 0.8-1cm 深绿色圆角矩形，"
                f"或其中间未包含白色电话图案"
            )

    return ok(
        "两页\"M.P.\"左侧均有 0.8-1cm 深绿色圆角矩形，中间含白色电话图案"
    )


def c13_email_text_both(model: PptModel) -> CheckResult:
    """+1：两页幻灯片均出现"z.xing@lanyunova.com"文本：从上至下位于页面
    70%-80% 高度，从左至右位于 10%-40% 宽度的区域范围内；字体为 Arial，
    13-17 磅。
    """
    for no in (1, 2):
        s = first_text(model, no, "z.xing@lanyunova.com")
        if not s:
            return fail(f"第{no}页未找到\"z.xing@lanyunova.com\"文本")

        # 位置："从上至下位于页面 70%-80% 高度"、"从左至右位于 10%-40% 宽度"。
        # 以形状左上角为基准（与本细则中其它同类"位于...区域范围内"条目
        # 保持一致的判定基准）。保留 2% 容差兼容办公软件 EMU→百分比的
        # 渲染精度损失。
        x_pct = s.x / model.slide_w
        y_pct = s.y / model.slide_h
        if not between(y_pct, 0.70, 0.80, tol=0.02):
            return fail(
                f"第{no}页邮箱纵向位置 {y_pct:.1%} 不在 70%-80%"
            )
        if not between(x_pct, 0.10, 0.40, tol=0.02):
            return fail(
                f"第{no}页邮箱横向位置 {x_pct:.1%} 不在 10%-40%"
            )

        # 字体：Arial。兼容 <a:latin>/<a:ea>/<a:cs> 三种 typeface 写入方式。
        try:
            typefaces = re.findall(r'typeface="([^"]+)"', s.raw._element.xml)
        except Exception:
            typefaces = []
        latin_names = font_names(s)
        all_names = [n for n in latin_names + typefaces if n]
        if not all_names or not all(is_arial(n) for n in all_names):
            return fail(
                f"第{no}页邮箱字体不是 Arial："
                f"latin={latin_names} typefaces={typefaces}"
            )

        # 字号：13-17 磅
        sizes = font_sizes(s)
        if not sizes or not all(between(sz, 13, 17) for sz in sizes):
            return fail(
                f"第{no}页邮箱字号不在 13-17 磅：{sizes}"
            )

    return ok(
        "两页\"z.xing@lanyunova.com\"位置、Arial 字体、13-17 磅均符合"
    )


def c14_email_icon_both(model: PptModel) -> CheckResult:
    """+5：两页幻灯片"z.xing@lanyunova.com"左侧均出现一个宽 0.8-1cm、高
    0.8-1cm 的深绿色填充圆角矩形形状，中间是一个白色信封图案。
    """
    envelope_glyphs = ["✉", "✉️", "envelope", "mail", "📧"]
    for no in (1, 2):
        email = first_text(model, no, "z.xing@lanyunova.com")
        if not email:
            return fail(
                f"第{no}页未找到\"z.xing@lanyunova.com\"，无法定位其左侧图标"
            )

        # 候选：位于邮箱左侧、宽高 0.8-1cm、深绿色填充的圆角矩形。
        # - 位于左侧：候选形状右边不超过邮箱左边（+0.15cm 办公软件文本框内边距）
        # - 圆角矩形：OOXML prst="roundRect"（严格四个圆角）
        # - 宽 0.8-1cm、高 0.8-1cm
        # - 深绿色填充：is_deep_green(fill)
        candidates: list[ShapeInfo] = []
        for s in model.slide(no):
            if s is email:
                continue
            if not _is_four_corner_rounded_rect(s):
                continue
            if not (between(s.w, 0.8, 1.0) and between(s.h, 0.8, 1.0)):
                continue
            if not is_deep_green(s.fill):
                continue
            if s.x2 > email.x + 0.15:
                continue
            candidates.append(s)

        # 中间是白色信封图案：
        # 路径 A：圆角矩形自身文本框包含信封字符（✉/envelope/mail 等），
        #         且字体颜色全为白色。
        # 路径 B：圆角矩形几何范围内存在一张图片（is_picture）— 视为
        #         "信封图案"的图片形式。
        # 路径 C：圆角矩形几何范围内存在白色矢量图形/线条组合 —— 覆盖
        #         "信封图案由白色 freeform / line / group 构成"这种常见写法。
        #         判定条件：
        #           a) 形状不是圆角矩形本身、不是邮箱文本；
        #           b) 形状不是图片（图片走路径 B）；
        #           c) 形状不承载文字（文字走路径 A）；
        #           d) 填充或线条颜色为白色；
        #           e) 几何上落入圆角矩形内部（面积≥50%重叠；线条按中心点判定）；
        #           f) 尺寸不超过圆角矩形自身（排除整页级白色背景）。
        def _mostly_inside(inner: "ShapeInfo", container: "ShapeInfo") -> bool:
            if inner.area > 0:
                return overlap_area(inner, container) / inner.area >= 0.50
            # 退化形状(w 或 h 为 0，如直线)：用中心点是否落入 container
            return (container.x - 0.05 <= inner.cx <= container.x2 + 0.05
                    and container.y - 0.05 <= inner.cy <= container.y2 + 0.05)

        def _is_white_vector_inside(inner: "ShapeInfo", container: "ShapeInfo", email_shape: "ShapeInfo") -> bool:
            if inner is container or inner is email_shape:
                return False
            if is_picture(inner):
                return False
            if inner.text.strip():
                return False
            if not (is_white(inner.fill) or is_white(inner.line)):
                return False
            # 排除超出圆角矩形自身尺寸的形状（页面级白色背景/大装饰）
            if inner.w > container.w * 1.1 or inner.h > container.h * 1.1:
                return False
            return _mostly_inside(inner, container)

        icon_ok = None
        for c in candidates:
            has_glyph = any(g in c.text for g in envelope_glyphs)
            colors = font_colors(c)
            white_glyph = has_glyph and colors and all(is_white(cc) for cc in colors)
            if white_glyph:
                icon_ok = c
                break

            has_pic_inside = any(
                is_picture(s)
                and overlap_area(s, c) / max(0.01, min(s.area, c.area)) >= 0.50
                for s in model.slide(no)
            )
            if has_pic_inside:
                icon_ok = c
                break

            # 路径 C：查找几何上落入圆角矩形内部的白色矢量图形/线条组合
            has_white_vector = any(
                _is_white_vector_inside(s, c, email) for s in model.slide(no)
            )
            if has_white_vector:
                icon_ok = c
                break

        if not icon_ok:
            return fail(
                f"第{no}页邮箱左侧未找到 0.8-1cm 深绿色圆角矩形，"
                f"或其中间未包含白色信封图案"
            )

    return ok("两页邮箱左侧均有 0.8-1cm 深绿色圆角矩形，中间含白色信封图案")


def c15_cn_address_text(model: PptModel) -> CheckResult:
    """+1：第1页幻灯片出现"福建省厦门市海沧区云创大道88号蓝湾创新中心12层
    1208室"文本：从上至下位于页面下方 80%-90% 高度，从左至右位于 10%-40%
    宽度的区域范围内；且在"z.xing@lanyunova.com"文本下方并与其左对齐；
    字体为微软雅黑或黑体 9-15 磅。
    """
    target = "福建省厦门市海沧区云创大道88号蓝湾创新中心12层1208室"
    addr = first_text(model, 1, target)
    email = first_text(model, 1, "z.xing@lanyunova.com")
    if not addr:
        return fail(f"第1页未找到中文地址文本：{target}")
    if not email:
        return fail("第1页未找到\"z.xing@lanyunova.com\"，无法判定地址与其左对齐")

    # 位置：从上至下位于 80%-90% 高度、从左至右位于 10%-40% 宽度。
    # 以形状左上角为基准（与本细则中其它同类"位于...区域范围内"条目
    # 保持一致）。保留 2% 容差兼容办公软件 EMU→百分比的渲染精度损失。
    x_pct = addr.x / model.slide_w
    y_pct = addr.y / model.slide_h
    if not between(y_pct, 0.80, 0.90, tol=0.02):
        return fail(f"第1页中文地址纵向位置 {y_pct:.1%} 不在 80%-90%")
    if not between(x_pct, 0.10, 0.40, tol=0.02):
        return fail(f"第1页中文地址横向位置 {x_pct:.1%} 不在 10%-40%")

    # 在邮箱下方：地址顶部在邮箱底部之下
    if not addr.y >= email.y2:
        return fail(
            f"第1页中文地址未位于邮箱下方："
            f"email_y2={email.y2:.2f} addr_y={addr.y:.2f}"
        )

    # 与邮箱左对齐：左边界差 ≤ 0.35cm（办公软件文本框内边距 ~0.25cm）
    if abs(addr.x - email.x) > 0.35:
        return fail(
            f"第1页中文地址与邮箱未左对齐：Δx={abs(addr.x - email.x):.2f}cm"
        )

    # 字体：微软雅黑或黑体。兼容 <a:latin>/<a:ea>/<a:cs> 三种 typeface 写入
    # 方式，以适配 PowerPoint / WPS / Keynote 对中文字体的不同存放位置。
    try:
        typefaces = re.findall(r'typeface="([^"]+)"', addr.raw._element.xml)
    except Exception:
        typefaces = []
    latin_names = font_names(addr)
    all_names = [n for n in latin_names + typefaces if n]
    if not all_names or not any(is_yahei_or_heiti(n) for n in all_names):
        return fail(
            f"第1页中文地址字体不是微软雅黑/黑体："
            f"latin={latin_names} typefaces={typefaces}"
        )

    # 字号：9-15 磅（每个非空 run 都必须落在区间内）
    sizes = font_sizes(addr)
    if not sizes or not all(between(sz, 9, 15) for sz in sizes):
        return fail(f"第1页中文地址字号不在 9-15 磅：{sizes}")

    return ok(
        f"第1页中文地址符合：pos=({x_pct:.1%},{y_pct:.1%}) 字号={sizes}"
    )


def c16_cn_address_icon(model: PptModel) -> CheckResult:
    target = "福建省厦门市海沧区云创大道88号蓝湾创新中心12层1208室"
    return icon_left_of_text(model, 1, target, ["⌖", "⌾", "📍", "location", "pin"])


def c17_vertical_line_both(model: PptModel) -> CheckResult:
    """+1：两页幻灯片从左至右页面 45%-55% 处出现一条绿色单实线：整体垂直
    与页面；长度 6-8.5cm。
    """
    for no in (1, 2):
        matched = []
        for s in model.slide(no):
            # 整体垂直于页面：宽 ≈ 0（OOXML 中垂线为 <p:sp> 或 <p:cxnSp> 且
            # ext.cx == 0；python-pptx 会把连接符宽度归一为 0 或极小值）。
            # 用 is_vertical_line：prst=="LINE" 或宽/高有一维为 0；高 ≥ 0.3cm、
            # 宽 ≤ 0.12cm，即"整体垂直"于页面。
            if not is_vertical_line(s):
                continue

            # 从左至右页面 45%-55% 处：整条线的横向位置（以线中心为基准，
            # 因为线的 x 与 x2 在归一化后接近同一值）。保留 2% 容差兼容
            # 办公软件 EMU→百分比的渲染精度损失。
            cx_pct = s.cx / model.slide_w
            if not between(cx_pct, 0.45, 0.55, tol=0.02):
                continue

            # 长度 6-8.5cm：竖线长度即高度
            if not between(s.h, 6, 8.5):
                continue

            # 绿色
            if not is_green(s.line):
                continue

            # 单实线：<a:ln> 内 <a:prstDash val="solid"> 或未声明；无 <a:custDash>
            if not _is_solid_line(s):
                continue

            matched.append(s)

        if not matched:
            return fail(
                f"第{no}页 45%-55% 宽度处未找到绿色单实线竖线（长度 6-8.5cm）"
            )
    return ok("两页均有 45%-55% 处、6-8.5cm 绿色单实线竖线")


def _is_rounded_rect_or_pill(shape: ShapeInfo) -> bool:
    """判定 shape 是否为"圆角矩形或胶囊形"形状。

    OOXML 中：
    - 四角圆角矩形 → prst="roundRect"（python-pptx: ROUNDED_RECTANGLE / ROUNDRECT）
    - 胶囊形（两端半圆的矩形）→ prst="pill"（较少见），或使用
      prst="roundRect" 并将圆角调节手柄 adj 拉满至 50000（即
      corner_radius == height/2，视觉即胶囊形）。python-pptx 均映射为
      ROUNDED_RECTANGLE。此外 PowerPoint 也可通过 prst="ellipse" 拉长
      成胶囊，但那不属于"矩形"类，此处不接受。
    """
    p = shape.prst.upper()
    return p in ("ROUNDED_RECTANGLE", "ROUNDRECT", "PILL")


def c18_cn_business_pill(model: PptModel) -> CheckResult:
    """+3：第1页幻灯片"福建蓝屿新材料科技有限公司"文本下方出现绿色填充的
    圆角矩形或胶囊形形状：从上至下位于页面 20%-30% 高度，从左至右位于
    50%-75% 宽度范围内；其内部出现"主营"文本：字体为微软雅黑、黑体，
    白色，14-16 磅。
    """
    company = first_text(model, 1, "福建蓝屿新材料科技有限公司")
    if not company:
        return fail("第1页未找到\"福建蓝屿新材料科技有限公司\"，无法定位其下方胶囊")

    # 收集绿色填充、位于 20%-30% × 50%-75% 区域内、位于公司名下方的"圆角
    # 矩形或胶囊形"形状。以形状左上角为基准判定位置。保留 2% 容差兼容
    # 办公软件 EMU→百分比的渲染精度损失。
    candidates = []
    for s in model.slide(1):
        if not _is_rounded_rect_or_pill(s):
            continue
        if not is_green(s.fill):
            continue
        x_pct = s.x / model.slide_w
        y_pct = s.y / model.slide_h
        if not between(y_pct, 0.20, 0.30, tol=0.02):
            continue
        if not between(x_pct, 0.50, 0.75, tol=0.02):
            continue
        if s.y < company.y2:
            continue
        candidates.append(s)

    if not candidates:
        return fail(
            "第1页公司名下方 20%-30%×50%-75% 未找到绿色填充圆角矩形/胶囊形"
        )

    # 其内部出现"主营"文本 + 字体 + 颜色 + 字号：
    # 优先匹配自身文本框内含"主营"的形状；否则支持外部"主营"文本框在几何
    # 上落入圆角矩形（重叠比例 ≥ 50%）。
    for c in candidates:
        # 路径 A：自身文本框包含"主营"
        pill_shape = c if text_contains(c, "主营") else None

        # 路径 B：外部"主营"文本框落入胶囊几何范围
        if pill_shape is None:
            for s in model.slide(1):
                if s is c:
                    continue
                if not text_contains(s, "主营"):
                    continue
                if c.area <= 0 or s.area <= 0:
                    continue
                if overlap_area(s, c) / max(0.01, min(s.area, c.area)) < 0.50:
                    continue
                pill_shape = s
                break

        if pill_shape is None:
            continue

        # 字体：微软雅黑、黑体（细则用"、"表示或）。兼容
        # <a:latin>/<a:ea>/<a:cs> 三种 typeface 写入方式。
        try:
            typefaces = re.findall(r'typeface="([^"]+)"', pill_shape.raw._element.xml)
        except Exception:
            typefaces = []
        latin_names = font_names(pill_shape)
        all_names = [n for n in latin_names + typefaces if n]
        if not all_names or not any(is_yahei_or_heiti(n) for n in all_names):
            return fail(
                f"第1页胶囊内\"主营\"字体不是微软雅黑/黑体："
                f"latin={latin_names} typefaces={typefaces}"
            )

        # 颜色：白色
        colors = font_colors(pill_shape)
        if not colors or not all(is_white(cc) for cc in colors):
            return fail(
                f"第1页胶囊内\"主营\"颜色不是白色：{[color_name(c) for c in colors]}"
            )

        # 字号：14-16 磅
        sizes = font_sizes(pill_shape)
        if not sizes or not all(between(sz, 14, 16) for sz in sizes):
            return fail(f"第1页胶囊内\"主营\"字号不在 14-16 磅：{sizes}")

        return ok(
            f"第1页绿色圆角矩形/胶囊 {c.w:.2f}×{c.h:.2f}cm 位于公司名下方，"
            f"内含\"主营\"白色微软雅黑/黑体 {sizes} 磅"
        )

    return fail("第1页公司名下方绿色圆角矩形/胶囊未包含符合规格的\"主营\"文本")


def check_business_list(model: PptModel, slide_no: int, headings: list[str], items: list[str], heading_font, item_font,
                        heading_size: tuple[float, float], item_size: tuple[float, float], heading_color, item_color) -> CheckResult:
    line_records = []
    for s in model.slide(slide_no):
        if not s.text.strip():
            continue
        for para in s.paragraphs:
            line_records.append((para["text"], s, para))
    cleaned_to_record = {clean_text(t): (t, s, p) for t, s, p in line_records}
    needed = headings + items
    missing = [t for t in needed if clean_text(t) not in cleaned_to_record]
    if missing:
        return fail(f"第{slide_no}页业务列表缺少：{missing}")

    # 每个文本单独一行：已按段落拆分，只要各目标文本能匹配单个段落即可。
    xs = []
    for t in needed:
        raw, s, para = cleaned_to_record[clean_text(t)]
        xs.append(round(s.x, 1))
        runs = [r for r in para["runs"] if r["text"].strip()]
        if not runs:
            return fail(f"第{slide_no}页 {raw} 无有效字体信息")
        if t in headings:
            for r in runs:
                if not (heading_font(r["font"] or "") and r["size"] is not None and between(float(r["size"]), *heading_size, tol=0.35) and heading_color(r["color"])):
                    return fail(f"第{slide_no}页标题 {raw} 字体/字号/颜色不满足")
        else:
            for r in runs:
                if not (item_font(r["font"] or "") and r["size"] is not None and between(float(r["size"]), *item_size, tol=0.35) and item_color(r["color"])):
                    return fail(f"第{slide_no}页条目 {raw} 字体/字号/颜色不满足")
    if max(xs) - min(xs) > 2.8:  # 标题和条目允许在同一左栏内不同缩进，但不应分散到多列
        return fail(f"第{slide_no}页业务列表未整体左对齐，x分布过宽：{xs}")
    return ok(f"第{slide_no}页业务列表目标文本均为独立段落，字体/字号/颜色和左对齐符合")


def _is_yahei(name: str) -> bool:
    """严格判定 typeface 是否为"微软雅黑"（不含"黑体"）。"""
    n = norm_font_name(name)
    return any(k in n for k in ["yahei", "microsoftyahei", "微软雅黑"])


def c19_cn_business_list(model: PptModel) -> CheckResult:
    """+3：第1页幻灯片"主营"文本下方出现"环保图层："、"水性树脂"、"功能助剂"、
    "高阻隔图层"、"绿色材料："、"可降解薄膜"、"纸塑复合材料"、"应用方案："、
    "工业胶粘剂"、"包装表面处理"这几个文本内容："环保图层：""绿色材料：""应用方案："
    文本颜色为绿色，字体为微软雅黑，13-15 磅；剩余文本颜色为黑色，字体为
    微软雅黑，11-13 磅；每一个文本单独一行，且全部左对齐。
    """
    headings = ["环保图层:", "绿色材料:", "应用方案:"]
    items = [
        "水性树脂", "功能助剂", "高阻隔图层", "可降解薄膜",
        "纸塑复合材料", "工业胶粘剂", "包装表面处理",
    ]
    needed = headings + items

    pill = first_text(model, 1, "主营")
    if not pill:
        return fail("第1页未找到\"主营\"，无法判定业务列表在其下方")

    # 收集所有非空段落（对应 OOXML `<a:p>`），保留段落所在形状 x 与段落
    # 层级的 XML 以便判定对齐属性。
    para_records = []  # (target, shape, para, para_xml)
    slide_shapes = model.slide(1)
    for s in slide_shapes:
        if not s.text.strip():
            continue
        try:
            para_nodes = re.findall(r"<a:p\b.*?</a:p>", s.raw._element.xml, re.DOTALL)
        except Exception:
            para_nodes = [""] * len(s.paragraphs)
        # 段落 XML 与 s.paragraphs 顺序对应（python-pptx 遍历顺序与 XML 一致）
        for idx, para in enumerate(s.paragraphs):
            para_xml = para_nodes[idx] if idx < len(para_nodes) else ""
            cleaned = clean_text(para["text"])
            for t in needed:
                if clean_text(t) == cleaned:
                    para_records.append((t, s, para, para_xml))
                    break

    matched_targets = {rec[0] for rec in para_records}
    missing = [t for t in needed if t not in matched_targets]
    if missing:
        return fail(f"第1页业务列表缺少文本：{missing}")

    # 每一个文本单独一行 → 已按 <a:p> 段落匹配，若同一段落被同时匹配到两个
    # 目标或段落文本远长于目标本身，视为未"单独一行"。
    for t, s, para, _ in para_records:
        if clean_text(para["text"]) != clean_text(t):
            return fail(f"第1页业务列表 {t!r} 所在段落包含其它内容：{para['text']!r}")

    # 位于"主营"文本下方
    for t, s, para, _ in para_records:
        if s.y < pill.y2 - 0.05:
            return fail(
                f"第1页业务列表 {t!r} 未位于\"主营\"下方："
                f"pill.y2={pill.y2:.2f} shape.y={s.y:.2f}"
            )

    # 全部左对齐：
    # 判据 1（几何）：所有段落所在形状的 x 相同 —— 允许 0.35cm 偏差（办公
    # 软件对文本框内边距/自动布局的合理误差）。
    xs = [round(rec[1].x, 2) for rec in para_records]
    if max(xs) - min(xs) > 0.35:
        return fail(f"第1页业务列表未左对齐（形状 x 分散过大）：{xs}")

    # 判据 2（段落属性）：如果段落 <a:pPr algn="..."/> 显式声明对齐，必须
    # 为 "l"（左对齐）；未声明视为左对齐（OOXML 默认值）。
    for t, _, _, para_xml in para_records:
        algn_m = re.search(r'<a:pPr[^>]*algn="([^"]+)"', para_xml)
        if algn_m and algn_m.group(1) != "l":
            return fail(f"第1页业务列表 {t!r} 段落对齐不是左对齐：algn={algn_m.group(1)}")

    # 逐条校验字体/字号/颜色
    # - 字体：微软雅黑（细则仅写"微软雅黑"）。兼容 <a:latin>/<a:ea>/<a:cs>。
    # - 颜色：标题绿；条目黑。
    # - 字号：标题 13-15 磅；条目 11-13 磅。
    for t, shape, para, para_xml in para_records:
        runs = [r for r in para["runs"] if r["text"].strip()]
        if not runs:
            return fail(f"第1页业务列表 {t!r} 无有效字体信息")

        # 段落级 typefaces 优先，未取到再从形状级 XML 兜底
        para_typefaces = re.findall(r'typeface="([^"]+)"', para_xml)
        for r in runs:
            name_candidates = [r["font"]] + para_typefaces
            name_candidates = [n for n in name_candidates if n]
            if not name_candidates or not all(_is_yahei(n) for n in name_candidates):
                return fail(
                    f"第1页业务列表 {t!r} 字体不是微软雅黑：{name_candidates}"
                )

            if r["size"] is None:
                return fail(f"第1页业务列表 {t!r} 未声明字号")

            if t in headings:
                if not between(float(r["size"]), 13, 15):
                    return fail(
                        f"第1页业务列表标题 {t!r} 字号不在 13-15 磅：{r['size']}"
                    )
                if not is_green(r["color"]):
                    return fail(
                        f"第1页业务列表标题 {t!r} 颜色不是绿色：{color_name(r['color'])}"
                    )
            else:
                if not between(float(r["size"]), 11, 13):
                    return fail(
                        f"第1页业务列表条目 {t!r} 字号不在 11-13 磅：{r['size']}"
                    )
                if not is_black(r["color"]):
                    return fail(
                        f"第1页业务列表条目 {t!r} 颜色不是黑色：{color_name(r['color'])}"
                    )

    return ok(
        "第1页业务列表 10 项文本齐全、单独成行、左对齐；"
        "标题绿色微软雅黑 13-15pt，条目黑色微软雅黑 11-13pt"
    )


def circle_icon_left(model: PptModel, slide_no: int, heading: str,
                     icon_color_check: Callable = is_green) -> CheckResult:
    h = first_text(model, slide_no, heading)
    if not h:
        return fail(f"第{slide_no}页未找到标题 {heading}")
    circles = [s for s in model.slide(slide_no) if s.x2 <= h.x + 0.25 and abs(s.cy - h.cy) <= 0.65 and is_oval(s) and between((s.w + s.h) / 2, 1.3, 1.5, 0.12) and is_light_gray(s.fill)]
    for c in circles:
        # 判定圆内"组合图形或图片"。
        # 1) 圆形本身是一张图片 → 直接通过
        if is_picture(c):
            return ok(f"第{slide_no}页 {heading} 左侧找到浅灰圆形图标（图片），直径约 {(c.w+c.h)/2:.2f}cm")
        # 2) 圆形几何范围内存在 ≥2 个独立"图形/图片"，且至少一个填充满足目标颜色
        #    → 视为"组合图形"。单个文本字符不计入"组合图形"。
        children = []
        target_color_hit = False
        for s in model.slide(slide_no):
            if s is c:
                continue
            inter = overlap_area(s, c)
            if min(c.area, max(s.area, 0.001)) <= 0:
                continue
            ratio = inter / max(0.01, min(s.area, c.area))
            if ratio < 0.40:
                continue
            # 圆形外框/同尺寸圆形不计入子元素
            if is_oval(s) and abs(s.w - c.w) < 0.1 and abs(s.h - c.h) < 0.1:
                continue
            if is_picture(s):
                children.append(s)
                target_color_hit = True
                continue
            # 几何形状（非纯文本框）：必须有可识别的填充或线条颜色
            has_geom = bool(s.prst) and s.prst not in ("RECTANGLE",) and (s.fill is not None or s.line is not None)
            if has_geom:
                children.append(s)
                if icon_color_check(s.fill) or icon_color_check(s.line):
                    target_color_hit = True
        if len(children) >= 2 and target_color_hit:
            return ok(f"第{slide_no}页 {heading} 左侧找到浅灰圆形图标（组合图形 {len(children)} 件），直径约 {(c.w+c.h)/2:.2f}cm")
    return fail(f"第{slide_no}页 {heading} 左侧未找到浅灰圆形及绿色主题组合图形/图片")


def c20_icon_eco_cn(model: PptModel) -> CheckResult:
    """+3：第1页幻灯片"环保涂层"文本左侧有一个圆形填充形状：直径在
    1.3-1.5cm，颜色为浅灰色；其中间有一个"绿色的水滴造型和一片小叶子"的
    组合图形或者整体是一个图片。
    """
    heading = first_text(model, 1, "环保涂层")
    if not heading:
        return fail("第1页未找到\"环保涂层\"文本，无法定位其左侧图标")

    slide = model.slide(1)

    # 候选圆形：
    # - 位于文本左侧：候选右边不超过文本左边 + 0.25cm（办公软件文本框内边距）
    # - 垂直中心与文本中心接近：|Δcy| ≤ 0.65cm
    # - 圆形：OOXML prst="ellipse"（python-pptx: OVAL/ELLIPSE）
    # - 直径 1.3-1.5cm：宽、高均需在区间内
    # - 浅灰色填充
    circles = []
    for s in slide:
        if s is heading:
            continue
        if s.x2 > heading.x + 0.25:
            continue
        if abs(s.cy - heading.cy) > 0.65:
            continue
        if not is_oval(s):
            continue
        if not (between(s.w, 1.3, 1.5) and between(s.h, 1.3, 1.5)):
            continue
        if not is_light_gray(s.fill):
            continue
        circles.append(s)

    if not circles:
        return fail(
            "第1页\"环保涂层\"左侧未找到浅灰色填充、直径 1.3-1.5cm 的圆形"
        )

    # 判定圆的"中间"是否有：
    # (A) 整体是一个图片：圆形本身为 <p:pic>（is_picture(circle) 为真）；
    #     或者圆形几何范围内包含一张图片（重叠比例 ≥ 50%）。
    # (B) 绿色水滴 + 小叶子的组合图形：需能静态区分"水滴"与"叶子"两个构件。
    #     - 水滴构件："prst" 名称含 TEAR/TEARDROP（PowerPoint 原生 teardrop
    #       预设几何），且填充或线条为绿色。
    #     - 叶子构件：与水滴不是同一件形状，且填充或线条为绿色。PowerPoint
    #       没有原生 leaf 预设，叶子绝大多数以 freeform (custGeom / 无 prst /
    #       FREEFORM 等) 方式绘制，静态无法从曲线路径判定"这画的是叶子"，
    #       故叶子放宽为"另一件绿色几何形状"，与水滴形成"至少两件、语义
    #       可区分"的组合。
    #     - 二者缺一 → 路径 B 不命中；只能走路径 A（图片）。
    #     组合还可能以 <p:grpSp>（GroupShape）形式存在——其子形状会由
    #     flatten 递归展开，计入 children。
    for c in circles:
        # (A-1) 圆形自身是图片
        if is_picture(c):
            return ok(
                f"第1页\"环保涂层\"左侧找到浅灰圆形图片，直径约 "
                f"{(c.w+c.h)/2:.2f}cm"
            )

        # 圆形几何内的其它形状（不含圆形自身、不含整页边框类超大形状）
        children: list[ShapeInfo] = []
        for s in slide:
            if s is c:
                continue
            # 排除整页级/接近整页的大形状（如边框圆角矩形）
            if s.area > c.area * 4:
                continue
            if min(s.area, c.area) <= 0:
                continue
            inter = overlap_area(s, c)
            if inter <= 0:
                continue
            if inter / max(0.01, min(s.area, c.area)) < 0.40:
                continue
            children.append(s)

        # (A-2) 圆内包含图片
        if any(is_picture(s) for s in children):
            return ok(
                f"第1页\"环保涂层\"左侧圆形中间包含图片，直径约 "
                f"{(c.w+c.h)/2:.2f}cm"
            )

        # (B) 组合图形：需分别识别"水滴"与"叶子"两件绿色构件
        combo: list[ShapeInfo] = [
            s for s in children
            if s.prst  # 需有 prstGeom 或图片标识
            and (s.fill is not None or s.line is not None or is_picture(s))
        ]
        # 水滴构件：prst 名称含 TEAR/TEARDROP，且颜色为绿色
        teardrops: list[ShapeInfo] = [
            s for s in combo
            if ("TEAR" in s.prst.upper() or "TEARDROP" in s.prst.upper())
            and (is_green(s.fill) or is_green(s.line))
        ]
        # 叶子构件：不是水滴，且颜色为绿色（叶子多为 freeform / custGeom，
        # 静态无法从曲线路径判定"这画的是叶子"，因此叶子按"另一件绿色几何
        # 构件"识别；只有当同时存在水滴构件时才算命中组合语义）。
        leaves: list[ShapeInfo] = [
            s for s in combo
            if s not in teardrops
            and (is_green(s.fill) or is_green(s.line))
        ]
        if teardrops and leaves:
            return ok(
                f"第1页\"环保涂层\"左侧圆形中间找到绿色水滴+叶子组合 "
                f"（水滴 {len(teardrops)} 件、叶子 {len(leaves)} 件）"
            )

    return fail(
        "第1页\"环保涂层\"左侧浅灰圆形中间未找到绿色水滴+叶子的组合图形或图片"
    )


def c21_icon_green_cn(model: PptModel) -> CheckResult:
    """+3：第1页幻灯片"绿色材料"文本左侧有一个圆形填充形状：直径在
    1.3-1.5cm，颜色为浅灰色；其中间有一个绿色的"带有白色脉络的叶子"组合
    图形或者整体是一个图片。
    """
    heading = first_text(model, 1, "绿色材料")
    if not heading:
        return fail("第1页未找到\"绿色材料\"文本，无法定位其左侧图标")

    slide = model.slide(1)

    # 候选圆形：位于文本左侧、直径 1.3-1.5cm、浅灰色填充
    circles = []
    for s in slide:
        if s is heading:
            continue
        if s.x2 > heading.x + 0.25:
            continue
        if abs(s.cy - heading.cy) > 0.65:
            continue
        if not is_oval(s):
            continue
        if not (between(s.w, 1.3, 1.5) and between(s.h, 1.3, 1.5)):
            continue
        if not is_light_gray(s.fill):
            continue
        circles.append(s)

    if not circles:
        return fail(
            "第1页\"绿色材料\"左侧未找到浅灰色填充、直径 1.3-1.5cm 的圆形"
        )

    # "其中间有一个绿色的带有白色脉络的叶子的组合图形或者整体是一个图片"
    # 分两条互斥路径：
    # 路径 A：整体是一个图片
    #   - A-1：圆形自身是 <p:pic>
    #   - A-2：圆形几何范围内包含一张图片（重叠比例 ≥ 50%）
    # 路径 B：组合图形（≥2 件独立几何形状/图片）
    #   - 位于圆形几何范围内（重叠比例 ≥ 40%，排除整页级超大形状）
    #   - 至少一件填充/线条为绿色（叶子主体）
    #   - 至少一件填充/线条为白色（脉络）
    for c in circles:
        # (A-1)
        if is_picture(c):
            return ok(
                f"第1页\"绿色材料\"左侧找到浅灰圆形图片，直径约 "
                f"{(c.w+c.h)/2:.2f}cm"
            )

        children = []
        for s in slide:
            if s is c:
                continue
            if s.area > c.area * 4:
                continue
            if min(s.area, c.area) <= 0:
                continue
            inter = overlap_area(s, c)
            if inter <= 0:
                continue
            if inter / max(0.01, min(s.area, c.area)) < 0.40:
                continue
            children.append(s)

        # (A-2)
        if any(is_picture(s) for s in children):
            return ok(
                f"第1页\"绿色材料\"左侧圆形中间包含图片，直径约 "
                f"{(c.w+c.h)/2:.2f}cm"
            )

        # (B) 组合图形：绿色叶子主体 + 白色脉络
        combo = [
            s for s in children
            if s.prst
            and (s.fill is not None or s.line is not None or is_picture(s))
        ]
        if len(combo) >= 2:
            has_green = any(is_green(s.fill) or is_green(s.line) for s in combo)
            has_white = any(is_white(s.fill) or is_white(s.line) for s in combo)
            if has_green and has_white:
                return ok(
                    f"第1页\"绿色材料\"左侧圆形中间找到组合图形（{len(combo)} 件），"
                    f"含绿色叶子与白色脉络"
                )

    return fail(
        "第1页\"绿色材料\"左侧浅灰圆形中间未找到绿色叶子+白色脉络的组合"
        "图形或图片"
    )


def c22_icon_app_cn(model: PptModel) -> CheckResult:
    """+3：第1页幻灯片"应用方案"文本左侧有一个圆形填充形状：直径在
    1.3-1.5cm，颜色为浅灰色；其中间有一个"绿色的带有白色箱盖和锁扣的立体
    纸箱"组合图型或者整体是一个图片。
    """
    heading = first_text(model, 1, "应用方案")
    if not heading:
        return fail("第1页未找到\"应用方案\"文本，无法定位其左侧图标")

    slide = model.slide(1)

    # 候选圆形：
    # - 位于文本左侧：候选右边不超过文本左边 + 0.25cm（办公软件文本框内边距）
    # - 垂直中心与文本中心接近：|Δcy| ≤ 0.65cm
    # - 圆形：OOXML prst="ellipse"（python-pptx: OVAL/ELLIPSE）
    # - 直径 1.3-1.5cm：宽、高均需在区间内
    # - 浅灰色填充（细则"颜色为浅灰色"→ fill 判定，兼容 PowerPoint/WPS/Keynote）
    circles = []
    for s in slide:
        if s is heading:
            continue
        if s.x2 > heading.x + 0.25:
            continue
        if abs(s.cy - heading.cy) > 0.65:
            continue
        if not is_oval(s):
            continue
        if not (between(s.w, 1.3, 1.5) and between(s.h, 1.3, 1.5)):
            continue
        if not is_light_gray(s.fill):
            continue
        circles.append(s)

    if not circles:
        return fail(
            "第1页\"应用方案\"左侧未找到浅灰色填充、直径 1.3-1.5cm 的圆形"
        )

    # "其中间有一个绿色的带有白色箱盖和锁扣的立体纸箱组合图型或者整体是
    # 一个图片"—— 分两条互斥路径：
    # 路径 A：整体是一个图片
    #   - A-1：圆形自身是 <p:pic>
    #   - A-2：圆形几何范围内包含一张图片（重叠比例 ≥ 50% 认为"在其中间"）
    # 路径 B：组合图形（≥2 件独立几何形状/图片）
    #   - 位于圆形几何范围内（重叠比例 ≥ 40%，排除整页级超大形状）
    #   - 至少一件填充/线条为绿色（纸箱主体）
    #   - 至少一件填充/线条为白色（箱盖 / 锁扣）
    for c in circles:
        # (A-1) 圆形自身是图片
        if is_picture(c):
            return ok(
                f"第1页\"应用方案\"左侧找到浅灰圆形图片，直径约 "
                f"{(c.w+c.h)/2:.2f}cm"
            )

        # 圆形几何内的其它形状（排除圆形自身与整页级超大形状）
        children = []
        for s in slide:
            if s is c:
                continue
            if s.area > c.area * 4:
                continue
            if min(s.area, c.area) <= 0:
                continue
            inter = overlap_area(s, c)
            if inter <= 0:
                continue
            if inter / max(0.01, min(s.area, c.area)) < 0.40:
                continue
            children.append(s)

        # (A-2) 圆形几何内包含一张图片
        for s in children:
            if is_picture(s) and overlap_area(s, c) / max(0.01, min(s.area, c.area)) >= 0.50:
                return ok(
                    f"第1页\"应用方案\"左侧圆形中间包含图片，直径约 "
                    f"{(c.w+c.h)/2:.2f}cm"
                )

        # (B) 组合图形：绿色纸箱主体 + 白色箱盖/锁扣
        combo = [
            s for s in children
            if s.prst
            and (s.fill is not None or s.line is not None or is_picture(s))
        ]
        if len(combo) >= 2:
            has_green = any(is_green(s.fill) or is_green(s.line) for s in combo)
            has_white = any(is_white(s.fill) or is_white(s.line) for s in combo)
            if has_green and has_white:
                return ok(
                    f"第1页\"应用方案\"左侧圆形中间找到组合图形（{len(combo)} 件），"
                    f"含绿色纸箱与白色箱盖/锁扣"
                )

    return fail(
        "第1页\"应用方案\"左侧浅灰圆形中间未找到绿色纸箱+白色箱盖/锁扣的"
        "组合图形或图片"
    )


def c23_slide2_title_separator(model: PptModel) -> CheckResult:
    """+3：第2页幻灯片"FUJIAN LANYU NOVA MATERIALS CO., LTD."文本左下方有
    两条直线：从上至下整体位于页面 20%-50% 处；两条线在同一水平线上；类型
    为单实线；颜色为橙色；长度 2.5-4cm；两条线中间有一个橙色或者黄色填充
    的叶片组合形状或者图片，整体高度 0.4-0.5cm，宽度 0.3-0.4cm。
    """
    title = first_text(model, 2, "FUJIAN LANYU NOVA MATERIALS CO., LTD.")
    if not title:
        return fail("第2页未找到英文公司名，无法定位其左下方分隔线")

    slide = model.slide(2)

    # 候选横线判定（细则逐点）：
    # A) "两条直线" → 水平线段：is_horizontal_line 判定 prst==LINE 或宽/高
    #    有一维为 0（办公软件写入 <p:cxnSp prstGeom prst="line"> 时 y、
    #    y+cy 相等，宽度 == 0 的直线段）
    # B) "从上至下整体位于页面 20%-50% 处" → 线的 y 与 y2 都需落入 20%-50%
    #    区间（保留 2% 容差以兼容办公软件 EMU→百分比的渲染精度损失）
    # C) "类型为单实线" → _is_solid_line：<a:ln> 内 <a:prstDash> 未声明
    #    或 val="solid"，且无 <a:custDash>（兼容 PowerPoint/WPS/Keynote）
    # D) "颜色为橙色" → is_orange(s.line)
    # E) "长度 2.5-4cm" → 水平线段的可视长度即 s.w
    # F) "文本左下方"（作用于整对，见后文 pair 判定，此处仅先过滤单线的
    #    位置约束——必须位于标题底部之下）
    lines = []
    for s in slide:
        if not is_horizontal_line(s):
            continue
        y1_pct = s.y / model.slide_h
        y2_pct = s.y2 / model.slide_h
        if not (between(y1_pct, 0.20, 0.50, tol=0.02)
                and between(y2_pct, 0.20, 0.50, tol=0.02)):
            continue
        if not between(s.w, 2.5, 4):
            continue
        if not is_orange(s.line):
            continue
        if not _is_solid_line(s):
            continue
        if s.y < title.y2:
            continue
        lines.append(s)

    # "两条线在同一水平线上"：cy 差 ≤ 0.10cm（≈同一渲染像素行）；且水平不
    # 重叠（`left.x2 ≤ right.x`），构成两条独立并排的横线。
    # "文本左下方"（对 pair 判定）：pair 整体位于标题左半侧 → 两线的横向
    # 中点 (left.x + right.x2) / 2 ≤ 标题水平中心 title.cx（保留 0.3cm 容差
    # 以兼容办公软件文本框内边距对可视中心的影响）。
    pair = None
    for i, a in enumerate(lines):
        for b in lines[i + 1:]:
            if abs(a.cy - b.cy) > 0.10:
                continue
            left, right = (a, b) if a.x < b.x else (b, a)
            if left.x2 > right.x:
                continue
            pair_mid_x = (left.x + right.x2) / 2
            if pair_mid_x > title.cx + 0.30:
                continue
            pair = (left, right)
            break
        if pair:
            break

    if not pair:
        return fail(
            "第2页英文公司名左下方 20%-50% 未找到两条在同一水平线、"
            "长度 2.5-4cm、橙色单实线的横线"
        )

    left, right = pair
    middle_x1, middle_x2 = left.x2, right.x
    middle_cy = (left.cy + right.cy) / 2

    # "两条线中间有一个橙色或者黄色填充的叶片组合形状或者图片，整体高度
    # 0.4-0.5cm，宽度 0.3-0.4cm"：
    # - "两条线中间"：候选形状的水平中心落在两线间水平段（放宽 0.15cm
    #   办公软件文本框边距容差），垂直中心与线中心接近（≤0.45cm）
    # - 分两条互斥路径（细则用"或者"，命中其一即可）：
    #   路径 A：图片（<p:pic>），宽 0.3-0.4cm、高 0.4-0.5cm
    #   路径 B：组合形状 —— 中间区域内 ≥2 件具备几何 prst 与填充/线条的
    #           图形；填充中至少存在"橙色或者黄色"（细则"或者"，命中其一
    #           即可）；外包围盒宽 0.3-0.4cm、高 0.4-0.5cm
    middle_shapes = [
        s for s in slide
        if s is not left and s is not right
        and middle_x1 - 0.15 <= s.cx <= middle_x2 + 0.15
        and abs(s.cy - middle_cy) <= 0.45
    ]

    # 路径 A：图片
    for pic in [s for s in middle_shapes if is_picture(s)]:
        if between(pic.w, 0.3, 0.4) and between(pic.h, 0.4, 0.5):
            return ok(
                f"第2页英文公司名左下方双橙线中间找到叶片图片 "
                f"{pic.w:.2f}×{pic.h:.2f}cm"
            )

    # 路径 B：组合形状
    combo = [
        s for s in middle_shapes
        if not is_picture(s)
        and s.prst  # 需有几何形状（排除纯文本框）
        and (s.fill is not None or s.line is not None)
    ]
    if len(combo) >= 2:
        has_orange_or_yellow = any(
            is_orange(s.fill) or is_yellow(s.fill) for s in combo
        )
        if has_orange_or_yellow:
            cx1 = min(s.x for s in combo)
            cx2 = max(s.x2 for s in combo)
            cy1 = min(s.y for s in combo)
            cy2 = max(s.y2 for s in combo)
            combo_w, combo_h = cx2 - cx1, cy2 - cy1
            if between(combo_w, 0.3, 0.4) and between(combo_h, 0.4, 0.5):
                return ok(
                    f"第2页英文公司名左下方双橙线中间找到叶片组合形状"
                    f"（{len(combo)} 件），整体 {combo_w:.2f}×{combo_h:.2f}cm，"
                    f"含橙色或黄色填充"
                )
            return fail(
                f"第2页英文公司名左下方双橙线中间叶片组合尺寸不满足："
                f"{combo_w:.2f}×{combo_h:.2f}cm"
            )
        return fail(
            f"第2页英文公司名左下方双橙线中间组合缺少橙色或黄色填充："
            f"{[color_name(s.fill) for s in combo]}"
        )

    return fail("第2页英文公司名左下方双橙线中间未找到叶片组合形状或图片")


def _is_ethan_script_font(name: str) -> bool:
    """严格判定 typeface 是否为细则列出的四种脚本字体之一：
    Brush Script MT / Script MT Bold / Alex Brush / Segoe Script。

    OOXML `<a:latin typeface="..."/>`（以及 <a:ea>/<a:cs>）会原样写入字体
    名称，办公软件（PowerPoint / WPS / Keynote）之间会保持一致。这里对
    名称做去空格、小写归一化后按子串匹配，同时兼容常见异体名：
      - "Brush Script MT" 也常写作 "BrushScriptMT" / "BrushScriptStd"
      - "Script MT Bold" 也常写作 "ScriptMTBold"
      - "Alex Brush" 也常写作 "AlexBrush" / "AlexBrush-Regular"
      - "Segoe Script" 也常写作 "SegoeScript"
    """
    n = norm_font_name(name)
    if not n:
        return False
    if "scriptmtbold" in n:
        return True
    if "brushscriptmt" in n or "brushscriptstd" in n:
        return True
    if "alexbrush" in n:
        return True
    if "segoescript" in n:
        return True
    return False


def c24_ethan(model: PptModel) -> CheckResult:
    """+1：第2页幻灯片出现"Ethan Shen"文本：位于页面从左至右 5%-45%、从上
    至下 30%-60% 范围内；字体为 Brush Script MT、Script MT Bold、
    Alex Brush、Segoe Script，25-35 磅。
    """
    s = first_text(model, 2, "Ethan Shen")
    if not s:
        return fail("第2页未找到\"Ethan Shen\"文本")

    # 位置：形状左上角为基准，位于页面从左至右 5%-45%、从上至下 30%-60%
    # 范围内。保留 2% 容差兼容办公软件 EMU→百分比的渲染精度损失。
    x_pct = s.x / model.slide_w
    y_pct = s.y / model.slide_h
    if not between(x_pct, 0.05, 0.45, tol=0.02):
        return fail(f"第2页\"Ethan Shen\"横向位置 {x_pct:.1%} 不在 5%-45%")
    if not between(y_pct, 0.30, 0.60, tol=0.02):
        return fail(f"第2页\"Ethan Shen\"纵向位置 {y_pct:.1%} 不在 30%-60%")

    # 字体：Brush Script MT / Script MT Bold / Alex Brush / Segoe Script。
    # 兼容 <a:latin>/<a:ea>/<a:cs> 三种 typeface 写入方式（Ethan Shen 为
    # 英文文本，PowerPoint 写入 <a:latin>；WPS/Keynote 有时也会写入
    # <a:ea>/<a:cs>）。所有非空 run 的所有 typeface 均需命中细则四种之一。
    try:
        typefaces = re.findall(r'typeface="([^"]+)"', s.raw._element.xml)
    except Exception:
        typefaces = []
    latin_names = font_names(s)
    all_names = [n for n in latin_names + typefaces if n]
    if not all_names or not all(_is_ethan_script_font(n) for n in all_names):
        return fail(
            f"第2页\"Ethan Shen\"字体不在 Brush Script MT / Script MT Bold "
            f"/ Alex Brush / Segoe Script 之中：latin={latin_names} "
            f"typefaces={typefaces}"
        )

    # 字号：25-35 磅（每个非空 run 都必须落在区间内）
    sizes = font_sizes(s)
    if not sizes or not all(between(sz, 25, 35) for sz in sizes):
        return fail(f"第2页\"Ethan Shen\"字号不在 25-35 磅：{sizes}")

    return ok(
        f"第2页\"Ethan Shen\"符合：pos=({x_pct:.1%},{y_pct:.1%}) "
        f"字号={sizes} 字体={all_names}"
    )


def c25_en_address_lines(model: PptModel) -> CheckResult:
    """+1：第2页幻灯片"z.xing@lanyunova.com"下方出现"Suite 1208, BlueBay
    Innovation Center,"、"No.88 Yunchuang Avenue,"、"Haicang District,
    Xiamen, Fujian, China"文本：每个文本单独位于一行；左对齐；字体为
    Arial，10-12 磅，黑色、绿色、灰色。
    """
    targets = [
        "Suite 1208, BlueBay Innovation Center,",
        "No.88 Yunchuang Avenue,",
        "Haicang District, Xiamen, Fujian, China",
    ]
    email = first_text(model, 2, "z.xing@lanyunova.com")
    if not email:
        return fail("第2页未找到\"z.xing@lanyunova.com\"，无法判定地址在其下方")

    # 收集第 2 页所有非空段落（对应 OOXML `<a:p>`）。三行地址可能位于
    # 同一个文本框中的三个段落，也可能各自独立文本框 —— 段落级判定同时
    # 覆盖两种办公软件常见写法（PowerPoint / WPS / Keynote）。
    para_records = []  # (target, shape, para, para_xml)
    for s in model.slide(2):
        if not s.text.strip():
            continue
        try:
            para_nodes = re.findall(r"<a:p\b.*?</a:p>", s.raw._element.xml, re.DOTALL)
        except Exception:
            para_nodes = [""] * len(s.paragraphs)
        for idx, para in enumerate(s.paragraphs):
            para_xml = para_nodes[idx] if idx < len(para_nodes) else ""
            cleaned = clean_text(para["text"])
            for t in targets:
                if clean_text(t) == cleaned:
                    para_records.append((t, s, para, para_xml))
                    break

    matched = {rec[0] for rec in para_records}
    missing = [t for t in targets if t not in matched]
    if missing:
        return fail(f"第2页英文地址缺少文本：{missing}")

    # "每个文本单独位于一行" → 目标文本必须等于段落全部内容（不含前后空白）
    for t, _, para, _ in para_records:
        if clean_text(para["text"]) != clean_text(t):
            return fail(
                f"第2页英文地址 {t!r} 所在段落包含其它内容：{para['text']!r}"
            )

    # 位于"z.xing@lanyunova.com"下方：段落所在形状顶端在邮箱底部之下
    # （允许 0.05cm 办公软件段落间距误差）
    for t, s, _, _ in para_records:
        if s.y < email.y2 - 0.05:
            return fail(
                f"第2页英文地址 {t!r} 未位于邮箱下方："
                f"email.y2={email.y2:.2f} shape.y={s.y:.2f}"
            )

    # "左对齐"：
    # 判据 1（几何）：三个段落所在形状的 x 相同 —— 允许 0.35cm 办公软件
    # 文本框内边距/自动布局误差。若三段位于同一文本框，形状 x 天然相同。
    xs = [round(rec[1].x, 2) for rec in para_records]
    if max(xs) - min(xs) > 0.35:
        return fail(f"第2页英文地址三行未左对齐（形状 x 分散过大）：{xs}")

    # 判据 2（段落属性）：段落 <a:pPr algn="..."/> 若显式声明，必须为
    # "l"（左对齐）；未声明视为左对齐（OOXML 默认值）。
    for t, _, _, para_xml in para_records:
        algn_m = re.search(r'<a:pPr[^>]*algn="([^"]+)"', para_xml)
        if algn_m and algn_m.group(1) != "l":
            return fail(
                f"第2页英文地址 {t!r} 段落对齐不是左对齐：algn={algn_m.group(1)}"
            )

    # 字体：Arial。兼容 <a:latin>/<a:ea>/<a:cs> 三种 typeface 写入方式；
    # 段落级 typefaces 优先，未取到再从形状级 XML 兜底。
    # 字号：10-12 磅（每个非空 run 都必须落在区间内）
    # 颜色：黑色、绿色、灰色（细则列举三种，命中其一即可）
    for t, shape, para, para_xml in para_records:
        runs = [r for r in para["runs"] if r["text"].strip()]
        if not runs:
            return fail(f"第2页英文地址 {t!r} 无有效字体信息")

        para_typefaces = re.findall(r'typeface="([^"]+)"', para_xml)
        for r in runs:
            name_candidates = [r["font"]] + para_typefaces
            name_candidates = [n for n in name_candidates if n]
            if not name_candidates or not all(is_arial(n) for n in name_candidates):
                return fail(
                    f"第2页英文地址 {t!r} 字体不是 Arial：{name_candidates}"
                )

            if r["size"] is None:
                return fail(f"第2页英文地址 {t!r} 未声明字号")
            if not between(float(r["size"]), 10, 12):
                return fail(
                    f"第2页英文地址 {t!r} 字号不在 10-12 磅：{r['size']}"
                )

            if not (is_black(r["color"]) or is_green(r["color"])
                    or is_gray(r["color"])):
                return fail(
                    f"第2页英文地址 {t!r} 颜色不是黑色/绿色/灰色："
                    f"{color_name(r['color'])}"
                )

    return ok(
        "第2页英文地址三行位于邮箱下方、单独成行、左对齐；"
        "Arial 10-12 磅，颜色为黑色/绿色/灰色"
    )


def c26_en_address_icon(model: PptModel) -> CheckResult:
    """+3：第2页幻灯片"Haicang District, Xiamen, Fujian, China"文本左侧有
    一个宽 0.8-1cm、高 0.8-1cm 的深绿色填充的圆角矩形形状，其中间是一个
    白色定位标记图案。
    """
    text = first_text(model, 2, "Haicang District, Xiamen, Fujian, China")
    if not text:
        return fail("第2页未找到英文地址末行")

    slide = model.slide(2)

    # 候选：位于文本左侧、宽/高 0.8-1cm、深绿色填充的圆角矩形。
    # - 位于左侧：候选形状右边不超过文本左边 + 0.15cm（办公软件文本框内边距）
    # - 圆角矩形：OOXML prst="roundRect"（严格四个圆角）
    #   [注意]细则用"圆角矩形"，这里严格采用 `_is_four_corner_rounded_rect`
    #   与 D2-12 / D2-14 / D2-16 保持一致；旧代码使用的 is_round_rect 会把
    #   1 圆角/2 圆角等变体也纳入，超出细则语义。
    # - 宽 0.8-1cm、高 0.8-1cm
    # - 深绿色填充：is_deep_green(fill)
    candidates: list[ShapeInfo] = []
    for s in slide:
        if s is text:
            continue
        if not _is_four_corner_rounded_rect(s):
            continue
        if not (between(s.w, 0.8, 1.0) and between(s.h, 0.8, 1.0)):
            continue
        if not is_deep_green(s.fill):
            continue
        if s.x2 > text.x + 0.15:
            continue
        candidates.append(s)

    # 中间是一个白色定位标记图案：
    # 路径 A：圆角矩形自身文本框内含定位字符（⌖/⌾/📍/location/pin），
    #         且字体颜色全为白色。这是最常见的图标写法（形状 + 文本符号）。
    # 路径 B：圆角矩形几何范围内存在一张图片（<p:pic>，重叠 ≥ 50%）—— 视
    #         为"定位图案"的图片形式；图片本身是否为定位图无法在 OOXML
    #         级别校验，采用几何包含 + "存在图片"作为可行判据。
    # 路径 C：圆角矩形几何范围内存在白色矢量定位标记组合 —— 覆盖"定位图案
    #         由白色 freeform / 圆 / 针脚线条 / group 构成"这种写法。
    #         定位标记的经典构成：外轮廓水滴形(teardrop)/椭圆(oval) + 内部
    #         针脚小圆(oval)；亦可能是任意 freeform 曲线。判定策略：
    #           a) 至少一件填充或线条为白色；
    #           b) 该形状不承载文字（文字走路径 A）、不是图片（图片走路径 B）；
    #           c) 几何上落入圆角矩形内部（面积≥50%重叠；线条按中心点判定）；
    #           d) 单件尺寸不超过圆角矩形自身（排除整页级白色背景）；
    #           e) 优先识别"外轮廓 + 内针脚"组合（外部为 teardrop / oval /
    #              freeform，内部包含小 oval 针脚且色调可分）；若不满足，退
    #              化为"至少两件白色矢量构件"（覆盖 freeform 线条组合写法）；
    #              再退化为"外轮廓为 teardrop/oval 的单件白色矢量图形"（覆盖
    #              整个定位标记就是一个白色水滴/椭圆的写法）。
    location_glyphs = ["⌖", "⌾", "📍", "location", "pin"]

    def _mostly_inside(inner: "ShapeInfo", container: "ShapeInfo") -> bool:
        if inner.area > 0:
            return overlap_area(inner, container) / inner.area >= 0.50
        # 退化形状(w 或 h 为 0，如直线)：用中心点是否落入 container
        return (container.x - 0.05 <= inner.cx <= container.x2 + 0.05
                and container.y - 0.05 <= inner.cy <= container.y2 + 0.05)

    def _is_white_vector_inside(inner: "ShapeInfo", container: "ShapeInfo") -> bool:
        if inner is container or inner is text:
            return False
        if is_picture(inner):
            return False
        if inner.text.strip():
            return False
        if not (is_white(inner.fill) or is_white(inner.line)):
            return False
        # 排除超出圆角矩形自身尺寸的形状（页面级白色背景/大装饰）
        if inner.w > container.w * 1.1 or inner.h > container.h * 1.1:
            return False
        return _mostly_inside(inner, container)

    for c in candidates:
        # 路径 A
        has_glyph = any(g in c.text for g in location_glyphs)
        colors = font_colors(c)
        white_glyph = has_glyph and colors and all(is_white(cc) for cc in colors)
        if white_glyph:
            return ok(
                "第2页英文地址末行左侧找到深绿色圆角矩形，中间含白色定位"
                "标记字符"
            )

        # 路径 B
        has_pic_inside = any(
            is_picture(s)
            and overlap_area(s, c) / max(0.01, min(s.area, c.area)) >= 0.50
            for s in slide
        )
        if has_pic_inside:
            return ok(
                "第2页英文地址末行左侧找到深绿色圆角矩形，中间含定位图片"
            )

        # 路径 C：白色矢量定位标记组合
        white_vec = [s for s in slide if _is_white_vector_inside(s, c)]
        if white_vec:
            # 优先识别"外轮廓 + 内针脚"结构：
            #   - 外轮廓：teardrop / oval / freeform / 无 prst 的 custGeom；
            #   - 内针脚：小圆(oval)，且几何上落入外轮廓包围盒内部。
            outers = [
                s for s in white_vec
                if ("TEAR" in s.prst.upper() or "TEARDROP" in s.prst.upper()
                    or "OVAL" in s.prst.upper() or "ELLIPSE" in s.prst.upper()
                    or "FREEFORM" in s.prst.upper() or s.prst == "")
            ]
            pins = [s for s in white_vec if is_oval(s)]
            outer_pin_hit = False
            for outer in outers:
                for pin in pins:
                    if pin is outer:
                        continue
                    if (pin.w < outer.w * 0.9 and pin.h < outer.h * 0.9
                            and _mostly_inside(pin, outer)):
                        outer_pin_hit = True
                        break
                if outer_pin_hit:
                    break
            if outer_pin_hit:
                return ok(
                    "第2页英文地址末行左侧找到深绿色圆角矩形，中间含白色定位标记矢量组合（外轮廓+内针脚）"
                )
            # 退化：多件白色矢量构件（覆盖 freeform 线条组合写法）
            if len(white_vec) >= 2:
                return ok(
                    f"第2页英文地址末行左侧找到深绿色圆角矩形，中间含白色定位标记矢量组合（{len(white_vec)} 件）"
                )
            # 再退化：单件白色 teardrop / oval / freeform 定位标记
            single = white_vec[0]
            if ("TEAR" in single.prst.upper()
                    or "TEARDROP" in single.prst.upper()
                    or "OVAL" in single.prst.upper()
                    or "ELLIPSE" in single.prst.upper()
                    or "FREEFORM" in single.prst.upper()):
                return ok(
                    "第2页英文地址末行左侧找到深绿色圆角矩形，中间含白色定位标记矢量图形"
                )

    return fail(
        "第2页英文地址末行左侧未找到 0.8-1cm 深绿色圆角矩形，或其中间未"
        "包含白色定位标记图案"
    )


def c27_main_business_pill(model: PptModel) -> CheckResult:
    """+3：第2页幻灯片"FUJIAN LANYU NOVA MATERIALS CO., LTD."右下方出现
    绿色填充的圆角矩形或胶囊形形状：从上至下位于页面 25%-40%，从左至右
    位于 50%-75% 的位置范围内；其内部出现"MAIN BUSINESS"文本：字体为
    Arial、黑体，白色，10-15 磅。
    """
    title = first_text(model, 2, "FUJIAN LANYU NOVA MATERIALS CO., LTD.")
    if not title:
        return fail("第2页未找到英文公司名，无法定位其右下方胶囊")

    slide = model.slide(2)

    # 收集绿色填充、位于 25%-40% × 50%-75% 区域内、位于公司名"右下方"的
    # "圆角矩形或胶囊形"形状。
    # - "圆角矩形或胶囊形"：OOXML `prst="roundRect"`（含 adj 拉满至胶囊视觉）
    #   或 `prst="pill"`（少见）；python-pptx 均映射为 ROUNDED_RECTANGLE
    # - 绿色填充：is_green(fill)
    # - 位置：形状左上角为基准，落入 25%-40% 纵向、50%-75% 横向（保留 2%
    #   容差兼容办公软件 EMU→百分比的渲染精度损失）
    # - 右下方：形状左端在标题水平中心之右（`s.x ≥ title.cx`，允许 0.3cm
    #   办公软件文本框内边距容差）；形状顶端在标题底部之下（`s.y ≥ title.y2`）
    candidates = []
    for s in slide:
        if not _is_rounded_rect_or_pill(s):
            continue
        if not is_green(s.fill):
            continue
        x_pct = s.x / model.slide_w
        y_pct = s.y / model.slide_h
        if not between(y_pct, 0.25, 0.40, tol=0.02):
            continue
        if not between(x_pct, 0.50, 0.75, tol=0.02):
            continue
        if s.y < title.y2:
            continue
        if s.x < title.cx - 0.30:
            continue
        candidates.append(s)

    if not candidates:
        return fail(
            "第2页英文公司名右下方 25%-40%×50%-75% 未找到绿色填充"
            "圆角矩形/胶囊形"
        )

    # 其内部出现"MAIN BUSINESS"文本，且字体/颜色/字号满足细则：
    # - 路径 A：胶囊自身文本框包含"MAIN BUSINESS"（最常见的办公软件写法）
    # - 路径 B：外部"MAIN BUSINESS"文本框在几何上落入胶囊范围（重叠 ≥ 50%）
    #
    # 字体：Arial、黑体（细则用"、"表示或；黑体覆盖 SimHei/微软雅黑等中文
    #   黑体族），命中其一即可。兼容 <a:latin>/<a:ea>/<a:cs> 三种 typeface。
    # 颜色：白色（每个非空 run 都为白色）
    # 字号：10-15 磅（每个非空 run）
    for c in candidates:
        # 路径 A
        pill_shape = c if text_contains(c, "MAIN BUSINESS") else None

        # 路径 B
        if pill_shape is None:
            for s in slide:
                if s is c:
                    continue
                if not text_contains(s, "MAIN BUSINESS"):
                    continue
                if c.area <= 0 or s.area <= 0:
                    continue
                if overlap_area(s, c) / max(0.01, min(s.area, c.area)) < 0.50:
                    continue
                pill_shape = s
                break

        if pill_shape is None:
            continue

        # 字体：Arial 或 黑体（细则原文"Arial、黑体"表示或）。此处"黑体"
        # 严格指 SimHei / Heiti，不包含微软雅黑等其它中文黑体族字体。
        try:
            typefaces = re.findall(r'typeface="([^"]+)"', pill_shape.raw._element.xml)
        except Exception:
            typefaces = []
        latin_names = font_names(pill_shape)
        all_names = [n for n in latin_names + typefaces if n]
        if not all_names:
            return fail(
                f"第2页胶囊内\"MAIN BUSINESS\"未声明字体："
                f"latin={latin_names} typefaces={typefaces}"
            )
        for n in all_names:
            if not (is_arial(n) or is_heiti_only(n)):
                return fail(
                    f"第2页胶囊内\"MAIN BUSINESS\"字体不是 Arial 或 黑体："
                    f"{n!r}（全部：latin={latin_names} typefaces={typefaces}）"
                )

        # 颜色：白色
        colors = font_colors(pill_shape)
        if not colors or not all(is_white(cc) for cc in colors):
            return fail(
                f"第2页胶囊内\"MAIN BUSINESS\"颜色不是白色："
                f"{[color_name(cc) for cc in colors]}"
            )

        # 字号：10-15 磅
        sizes = font_sizes(pill_shape)
        if not sizes or not all(between(sz, 10, 15) for sz in sizes):
            return fail(
                f"第2页胶囊内\"MAIN BUSINESS\"字号不在 10-15 磅：{sizes}"
            )

        return ok(
            f"第2页绿色圆角矩形/胶囊 {c.w:.2f}×{c.h:.2f}cm 位于公司名右下方，"
            f"内含\"MAIN BUSINESS\"白色 Arial/黑体 {sizes} 磅"
        )

    return fail(
        "第2页公司名右下方绿色圆角矩形/胶囊未包含符合规格的"
        "\"MAIN BUSINESS\"文本"
    )


def c28_en_business_list(model: PptModel) -> CheckResult:
    """+3：第2页幻灯片"MAIN BUSINESS"下方出现"ECO COATINGS:"、
    "Water-based resins"、"Functional additives"、"Barrier coatings"、
    "GREEN MATERIALS:"、"Biodegradable films"、"Paper-plastic composites"、
    "APPLICATION SOLUTIONS:"、"Industrial adhesives"、
    "Packaging surface treatment"文本：每一个文本都单独排列在一行，全部
    左对齐；"ECO COATINGS:"、"GREEN MATERIALS:"、"APPLICATION SOLUTIONS:"
    这三个文本字体为绿色，Arial，11-13 磅；其它文本字体为黑色、灰色、
    绿色，Arial，9-11 磅。
    """
    headings = ["ECO COATINGS:", "GREEN MATERIALS:", "APPLICATION SOLUTIONS:"]
    items = [
        "Water-based resins", "Functional additives", "Barrier coatings",
        "Biodegradable films", "Paper-plastic composites",
        "Industrial adhesives", "Packaging surface treatment",
    ]
    needed = headings + items

    pill = first_text(model, 2, "MAIN BUSINESS")
    if not pill:
        return fail("第2页未找到\"MAIN BUSINESS\"，无法判定业务列表在其下方")

    # 收集所有非空段落（对应 OOXML `<a:p>`）。10 项文本可能位于同一个文本
    # 框中的多个段落，也可能各自独立文本框 —— 段落级判定同时覆盖两种办公
    # 软件常见写法（PowerPoint / WPS / Keynote）。
    para_records = []  # (target, shape, para, para_xml)
    for s in model.slide(2):
        if not s.text.strip():
            continue
        try:
            para_nodes = re.findall(r"<a:p\b.*?</a:p>", s.raw._element.xml, re.DOTALL)
        except Exception:
            para_nodes = [""] * len(s.paragraphs)
        for idx, para in enumerate(s.paragraphs):
            para_xml = para_nodes[idx] if idx < len(para_nodes) else ""
            cleaned = clean_text(para["text"])
            for t in needed:
                if clean_text(t) == cleaned:
                    para_records.append((t, s, para, para_xml))
                    break

    matched = {rec[0] for rec in para_records}
    missing = [t for t in needed if t not in matched]
    if missing:
        return fail(f"第2页英文业务列表缺少文本：{missing}")

    # "每一个文本都单独排列在一行" → 段落文本 clean_text 必须等于目标文本，
    # 不能包含其它内容（同段落多目标或段落含额外文本均视为不满足）
    for t, _, para, _ in para_records:
        if clean_text(para["text"]) != clean_text(t):
            return fail(
                f"第2页英文业务列表 {t!r} 所在段落包含其它内容：{para['text']!r}"
            )

    # 位于"MAIN BUSINESS"下方（细则"下方"）：所在形状顶端在 MAIN BUSINESS
    # 底部之下（允许 0.05cm 办公软件段落间距误差）
    for t, s, _, _ in para_records:
        if s.y < pill.y2 - 0.05:
            return fail(
                f"第2页英文业务列表 {t!r} 未位于\"MAIN BUSINESS\"下方："
                f"pill.y2={pill.y2:.2f} shape.y={s.y:.2f}"
            )

    # "全部左对齐"：
    # 判据 1（几何）：所有段落所在形状的 x 相同 —— 允许 0.35cm 办公软件
    #   文本框内边距/自动布局误差。同一文本框内多段天然共享 x。
    # 判据 2（段落属性）：段落 <a:pPr algn="..."/> 若显式声明，必须为
    #   "l"（左对齐）；未声明视为左对齐（OOXML 默认值）。
    xs = [round(rec[1].x, 2) for rec in para_records]
    if max(xs) - min(xs) > 0.35:
        return fail(f"第2页英文业务列表未左对齐（形状 x 分散过大）：{xs}")

    for t, _, _, para_xml in para_records:
        algn_m = re.search(r'<a:pPr[^>]*algn="([^"]+)"', para_xml)
        if algn_m and algn_m.group(1) != "l":
            return fail(
                f"第2页英文业务列表 {t!r} 段落对齐不是左对齐："
                f"algn={algn_m.group(1)}"
            )

    # 逐条校验字体/字号/颜色：
    # - 三个标题：Arial，11-13 磅，绿色
    # - 其它 7 条：Arial，9-11 磅，黑色/灰色/绿色（"、"表示或，命中其一即可）
    # 字体兼容 <a:latin>/<a:ea>/<a:cs> 三种 typeface 写入方式（段落级优先，
    # 形状级兜底）；颜色/字号逐 run 检查。
    for t, shape, para, para_xml in para_records:
        runs = [r for r in para["runs"] if r["text"].strip()]
        if not runs:
            return fail(f"第2页英文业务列表 {t!r} 无有效字体信息")

        para_typefaces = re.findall(r'typeface="([^"]+)"', para_xml)
        for r in runs:
            name_candidates = [r["font"]] + para_typefaces
            name_candidates = [n for n in name_candidates if n]
            if not name_candidates or not all(is_arial(n) for n in name_candidates):
                return fail(
                    f"第2页英文业务列表 {t!r} 字体不是 Arial：{name_candidates}"
                )

            if r["size"] is None:
                return fail(f"第2页英文业务列表 {t!r} 未声明字号")

            if t in headings:
                if not between(float(r["size"]), 11, 13):
                    return fail(
                        f"第2页英文业务列表标题 {t!r} 字号不在 11-13 磅："
                        f"{r['size']}"
                    )
                if not is_green(r["color"]):
                    return fail(
                        f"第2页英文业务列表标题 {t!r} 颜色不是绿色："
                        f"{color_name(r['color'])}"
                    )
            else:
                if not between(float(r["size"]), 9, 11):
                    return fail(
                        f"第2页英文业务列表条目 {t!r} 字号不在 9-11 磅："
                        f"{r['size']}"
                    )
                if not (is_black(r["color"]) or is_gray(r["color"])
                        or is_green(r["color"])):
                    return fail(
                        f"第2页英文业务列表条目 {t!r} 颜色不是黑色/灰色/绿色："
                        f"{color_name(r['color'])}"
                    )

    return ok(
        "第2页英文业务列表 10 项文本齐全、单独成行、左对齐；"
        "标题绿色 Arial 11-13 磅，条目 Arial 9-11 磅（黑/灰/绿色）"
    )


def c29_icon_eco_en(model: PptModel) -> CheckResult:
    """+3：第2页幻灯片"ECO COATINGS"左侧有一个圆形填充形状：直径在
    1.3-1.5cm，颜色为浅灰色；其中间有一个"绿色的水滴造型和一片小叶子"的
    组合图形或者此部分直接为一张图片。
    """
    heading = first_text(model, 2, "ECO COATINGS")
    if not heading:
        return fail("第2页未找到\"ECO COATINGS\"文本，无法定位其左侧图标")

    slide = model.slide(2)

    # 候选圆形：
    # - 位于文本左侧：候选右边不超过文本左边 + 0.25cm（办公软件文本框内边距）
    # - 垂直中心与文本中心接近：|Δcy| ≤ 0.65cm
    # - 圆形：OOXML prst="ellipse"（python-pptx: OVAL/ELLIPSE）
    # - 直径 1.3-1.5cm：宽、高均需在区间内
    # - 浅灰色填充
    circles = []
    for s in slide:
        if s is heading:
            continue
        if s.x2 > heading.x + 0.25:
            continue
        if abs(s.cy - heading.cy) > 0.65:
            continue
        if not is_oval(s):
            continue
        if not (between(s.w, 1.3, 1.5) and between(s.h, 1.3, 1.5)):
            continue
        if not is_light_gray(s.fill):
            continue
        circles.append(s)

    if not circles:
        return fail(
            "第2页\"ECO COATINGS\"左侧未找到浅灰色填充、直径 1.3-1.5cm 的圆形"
        )

    # 其中间有一个"绿色的水滴造型和一片小叶子"的组合图形，或者此部分直接
    # 为一张图片。分两条互斥路径：
    # 路径 A：此部分直接为一张图片
    #   - A-1：圆形自身是 <p:pic>
    #   - A-2：圆形几何范围内包含一张图片（重叠比例 ≥ 50% 视为"在其中间"）
    # 路径 B：组合图形——必须同时识别到「绿色水滴」与「绿色叶片」两个语义
    #   构件，才算命中，避免"只要圆内有 ≥2 件几何且有一件绿色"的宽松判定：
    #   - 位于圆形几何范围内（重叠比例 ≥ 40%，排除整页级超大形状）；
    #   - 水滴：prst 含 "TEAR"/"TEARDROP"，填充或线条为绿色；
    #   - 叶片：绿色的非水滴几何形状（可为 leaf/oval/freeform/blob 等）。
    for c in circles:
        # (A-1) 圆形自身是图片
        if is_picture(c):
            return ok(
                f"第2页\"ECO COATINGS\"左侧找到浅灰圆形图片，直径约 "
                f"{(c.w+c.h)/2:.2f}cm"
            )

        # 圆形几何内的其它形状（排除圆形自身与整页级超大形状）
        children: list[ShapeInfo] = []
        for s in slide:
            if s is c:
                continue
            if s.area > c.area * 4:
                continue
            if min(s.area, c.area) <= 0:
                continue
            inter = overlap_area(s, c)
            if inter <= 0:
                continue
            if inter / max(0.01, min(s.area, c.area)) < 0.40:
                continue
            children.append(s)

        # (A-2) 圆形几何内包含一张图片
        for s in children:
            if is_picture(s) and overlap_area(s, c) / max(0.01, min(s.area, c.area)) >= 0.50:
                return ok(
                    f"第2页\"ECO COATINGS\"左侧圆形中间包含图片，直径约 "
                    f"{(c.w+c.h)/2:.2f}cm"
                )

        # (B) 绿色水滴 + 绿色叶片
        combo: list[ShapeInfo] = [
            s for s in children
            if s.prst
            and (s.fill is not None or s.line is not None or is_picture(s))
        ]
        teardrops: list[ShapeInfo] = [
            s for s in combo
            if ("TEAR" in s.prst.upper() or "TEARDROP" in s.prst.upper())
            and (is_green(s.fill) or is_green(s.line))
        ]
        leaves: list[ShapeInfo] = [
            s for s in combo
            if s not in teardrops
            and (is_green(s.fill) or is_green(s.line))
        ]
        if teardrops and leaves:
            return ok(
                f"第2页\"ECO COATINGS\"左侧圆形中间找到绿色水滴+叶子组合（水滴 {len(teardrops)} 件、叶子 {len(leaves)} 件）"
            )

    return fail(
        "第2页\"ECO COATINGS\"左侧浅灰圆形中间未找到绿色水滴+小叶子的组合图形或图片"
    )


def c30_icon_green_en(model: PptModel) -> CheckResult:
    """+3：第2页幻灯片"GREEN MATERIALS"左侧有一个圆形填充形状：直径在
    1.3-1.5cm，颜色为浅灰色；其中间有一个绿色的"带有白色脉络的叶子"组合
    图形或者整体是一个图片。
    """
    heading = first_text(model, 2, "GREEN MATERIALS")
    if not heading:
        return fail("第2页未找到\"GREEN MATERIALS\"文本，无法定位其左侧图标")

    slide = model.slide(2)

    # 候选圆形：
    # - 位于文本左侧：候选右边不超过文本左边 + 0.25cm（办公软件文本框内边距）
    # - 垂直中心与文本中心接近：|Δcy| ≤ 0.65cm
    # - 圆形：OOXML prst="ellipse"（python-pptx: OVAL/ELLIPSE）
    # - 直径 1.3-1.5cm：宽、高均需在区间内
    # - 浅灰色填充
    circles = []
    for s in slide:
        if s is heading:
            continue
        if s.x2 > heading.x + 0.25:
            continue
        if abs(s.cy - heading.cy) > 0.65:
            continue
        if not is_oval(s):
            continue
        if not (between(s.w, 1.3, 1.5) and between(s.h, 1.3, 1.5)):
            continue
        if not is_light_gray(s.fill):
            continue
        circles.append(s)

    if not circles:
        return fail(
            "第2页\"GREEN MATERIALS\"左侧未找到浅灰色填充、直径 1.3-1.5cm 的圆形"
        )

    # 其中间有一个绿色的"带有白色脉络的叶子"组合图形，或者整体是一个图片。
    # 分两条互斥路径：
    # 路径 A：整体是一个图片
    #   - A-1：圆形自身是 <p:pic>
    #   - A-2：圆形几何范围内包含一张图片（重叠比例 ≥ 50% 视为"在其中间"）
    # 路径 B：组合图形（≥2 件独立几何形状/图片）
    #   - 位于圆形几何范围内（重叠比例 ≥ 40%，排除整页级超大形状）
    #   - 至少一件填充/线条为绿色（叶子主体）
    #   - 至少一件填充/线条为白色（脉络）
    for c in circles:
        # (A-1) 圆形自身是图片
        if is_picture(c):
            return ok(
                f"第2页\"GREEN MATERIALS\"左侧找到浅灰圆形图片，直径约 "
                f"{(c.w+c.h)/2:.2f}cm"
            )

        # 圆形几何内的其它形状（排除圆形自身与整页级超大形状）
        children = []
        for s in slide:
            if s is c:
                continue
            if s.area > c.area * 4:
                continue
            if min(s.area, c.area) <= 0:
                continue
            inter = overlap_area(s, c)
            if inter <= 0:
                continue
            if inter / max(0.01, min(s.area, c.area)) < 0.40:
                continue
            children.append(s)

        # (A-2) 圆形几何内包含一张图片
        for s in children:
            if is_picture(s) and overlap_area(s, c) / max(0.01, min(s.area, c.area)) >= 0.50:
                return ok(
                    f"第2页\"GREEN MATERIALS\"左侧圆形中间包含图片，直径约 "
                    f"{(c.w+c.h)/2:.2f}cm"
                )

        # (B) 组合图形：绿色叶子主体 + 白色脉络
        combo = [
            s for s in children
            if s.prst
            and (s.fill is not None or s.line is not None or is_picture(s))
        ]
        if len(combo) >= 2:
            has_green = any(is_green(s.fill) or is_green(s.line) for s in combo)
            has_white = any(is_white(s.fill) or is_white(s.line) for s in combo)
            if has_green and has_white:
                return ok(
                    f"第2页\"GREEN MATERIALS\"左侧圆形中间找到组合图形"
                    f"（{len(combo)} 件），含绿色叶子与白色脉络"
                )

    return fail(
        "第2页\"GREEN MATERIALS\"左侧浅灰圆形中间未找到绿色叶子+白色脉络"
        "的组合图形或图片"
    )


def c31_icon_app_en(model: PptModel) -> CheckResult:
    """+3：第2页幻灯片"APPLICATION SOLUTIONS"左侧有一个圆形填充形状：直径
    在 1.3-1.5cm，颜色为浅灰色；其中间有一个"绿色的带有白色箱盖和锁扣的
    立体纸箱"组合图型或者整体是一个图片。
    """
    heading = first_text(model, 2, "APPLICATION SOLUTIONS")
    if not heading:
        return fail("第2页未找到\"APPLICATION SOLUTIONS\"文本，无法定位其左侧图标")

    slide = model.slide(2)

    # 候选圆形：
    # - 位于文本左侧：候选右边不超过文本左边 + 0.25cm（办公软件文本框内边距）
    # - 垂直中心与文本中心接近：|Δcy| ≤ 0.65cm
    # - 圆形：OOXML prst="ellipse"（python-pptx: OVAL/ELLIPSE）
    # - 直径 1.3-1.5cm：宽、高均需在区间内
    # - 浅灰色填充
    circles = []
    for s in slide:
        if s is heading:
            continue
        if s.x2 > heading.x + 0.25:
            continue
        if abs(s.cy - heading.cy) > 0.65:
            continue
        if not is_oval(s):
            continue
        if not (between(s.w, 1.3, 1.5) and between(s.h, 1.3, 1.5)):
            continue
        if not is_light_gray(s.fill):
            continue
        circles.append(s)

    if not circles:
        return fail(
            "第2页\"APPLICATION SOLUTIONS\"左侧未找到浅灰色填充、"
            "直径 1.3-1.5cm 的圆形"
        )

    # "其中间有一个绿色的带有白色箱盖和锁扣的立体纸箱组合图型或者整体是
    # 一个图片"—— 分两条互斥路径：
    # 路径 A：整体是一个图片
    #   - A-1：圆形自身是 <p:pic>
    #   - A-2：圆形几何范围内包含一张图片（重叠比例 ≥ 50% 视为"在其中间"）
    # 路径 B：组合图形（≥2 件独立几何形状/图片）
    #   - 位于圆形几何范围内（重叠比例 ≥ 40%，排除整页级超大形状）
    #   - 至少一件填充/线条为绿色（纸箱主体）
    #   - 至少一件填充/线条为白色（箱盖 / 锁扣）
    for c in circles:
        # (A-1) 圆形自身是图片
        if is_picture(c):
            return ok(
                f"第2页\"APPLICATION SOLUTIONS\"左侧找到浅灰圆形图片，"
                f"直径约 {(c.w+c.h)/2:.2f}cm"
            )

        # 圆形几何内的其它形状（排除圆形自身与整页级超大形状）
        children = []
        for s in slide:
            if s is c:
                continue
            if s.area > c.area * 4:
                continue
            if min(s.area, c.area) <= 0:
                continue
            inter = overlap_area(s, c)
            if inter <= 0:
                continue
            if inter / max(0.01, min(s.area, c.area)) < 0.40:
                continue
            children.append(s)

        # (A-2) 圆形几何内包含一张图片
        for s in children:
            if is_picture(s) and overlap_area(s, c) / max(0.01, min(s.area, c.area)) >= 0.50:
                return ok(
                    f"第2页\"APPLICATION SOLUTIONS\"左侧圆形中间包含图片，"
                    f"直径约 {(c.w+c.h)/2:.2f}cm"
                )

        # (B) 组合图形：绿色纸箱主体 + 白色箱盖/锁扣
        combo = [
            s for s in children
            if s.prst
            and (s.fill is not None or s.line is not None or is_picture(s))
        ]
        if len(combo) >= 2:
            has_green = any(is_green(s.fill) or is_green(s.line) for s in combo)
            has_white = any(is_white(s.fill) or is_white(s.line) for s in combo)
            if has_green and has_white:
                return ok(
                    f"第2页\"APPLICATION SOLUTIONS\"左侧圆形中间找到组合"
                    f"图形（{len(combo)} 件），含绿色纸箱与白色箱盖/锁扣"
                )

    return fail(
        "第2页\"APPLICATION SOLUTIONS\"左侧浅灰圆形中间未找到绿色纸箱+"
        "白色箱盖/锁扣的组合图形或图片"
    )


def c32_bottom_arcs_slide1(model: PptModel) -> CheckResult:
    """+3：第1页幻灯片从下至上页面底部 10% 高度范围内有绿色和橙色的弧形
    色块、橙色面积小于绿色。
    """
    slide = model.slide(1)

    # "弧形色块"：OOXML `prstGeom` 中含有弧线 / 曲线的几何形状。常见的
    # 底部装饰"弧形色块"通常写作：
    #   - <a:prstGeom prst="wave"/>、prst="doubleWave"（波浪形）
    #   - prst="arc"（弧线）
    #   - prst="chord"（弓形）
    #   - prst="pie"（饼形 / 扇形）
    #   - prst="ellipse"（椭圆）
    #   - prst="moon"（月牙）
    #   - prst="teardrop"（泪滴）
    #   - prst="cloud"（云）
    # 上述形状在 PowerPoint / WPS / Keynote 中的写法一致，通过 prst 名称
    # 即可判定，跨办公软件有效。
    # 严格排除直线多边形（如 trapezoid / rectangle / triangle 等）——细则
    # 明确要求"弧形色块"。
    ARC_KEYWORDS = (
        "WAVE", "DOUBLE_WAVE", "ARC", "CHORD", "PIE",
        "ELLIPSE", "OVAL", "MOON", "TEAR", "TEARDROP", "CLOUD",
    )

    def is_arc_like(shape) -> bool:
        p = shape.prst.upper()
        return any(k in p for k in ARC_KEYWORDS)

    # 位置："从下至上页面底部 10% 高度范围内" → 形状整体位于页面纵向
    # 90%-100% 区间内。以形状顶部 s.y 落入区间为基准 —— 顶部在 90% 以下
    # 即整体在底部 10% 内。保留 2% 容差兼容办公软件 EMU→百分比的渲染精度
    # 损失。
    bottom_arcs = [
        s for s in slide
        if is_arc_like(s)
        and s.y / model.slide_h >= 0.90 - 0.02
    ]

    # 分别统计绿色和橙色弧形色块的填充面积（以外包围盒面积近似——
    # OOXML 层面无法访问曲线积分面积，各家办公软件亦以包围盒近似渲染）。
    green_shapes = [s for s in bottom_arcs if is_green(s.fill)]
    orange_shapes = [s for s in bottom_arcs if is_orange(s.fill)]
    green_area = sum(s.area for s in green_shapes)
    orange_area = sum(s.area for s in orange_shapes)

    # 有绿色弧形色块
    if green_area <= 0:
        return fail("第1页底部 10% 未找到绿色弧形色块")
    # 有橙色弧形色块
    if orange_area <= 0:
        return fail("第1页底部 10% 未找到橙色弧形色块")
    # 橙色面积小于绿色
    if orange_area >= green_area:
        return fail(
            f"第1页底部 10% 橙色弧形色块面积 {orange_area:.2f}cm² "
            f"未小于绿色 {green_area:.2f}cm²"
        )

    return ok(
        f"第1页底部 10% 检测到绿色/橙色弧形色块，"
        f"橙色面积 {orange_area:.2f}cm² 小于绿色 {green_area:.2f}cm²"
    )


def c37_bottom_arcs_slide2(model: PptModel) -> CheckResult:
    """+3：第2页幻灯片从下至上页面底部 20% 高度范围内有绿色和橙色的弧形
    色块、橙色面积小于绿色。
    """
    slide = model.slide(2)

    # "弧形色块"：OOXML `prstGeom` 中含有弧线 / 曲线的几何形状。常见的
    # 底部装饰"弧形色块"通常写作：
    #   - <a:prstGeom prst="wave"/>、prst="doubleWave"（波浪形）
    #   - prst="arc"（弧线）
    #   - prst="chord"（弓形）
    #   - prst="pie"（饼形 / 扇形）
    #   - prst="ellipse"（椭圆）
    #   - prst="moon"（月牙）
    #   - prst="teardrop"（泪滴）
    #   - prst="cloud"（云）
    # 上述形状在 PowerPoint / WPS / Keynote 中的写法一致，通过 prst 名称
    # 即可判定，跨办公软件有效。
    # 严格排除直线多边形（如 trapezoid / rectangle / triangle 等）——细则
    # 明确要求"弧形色块"。
    ARC_KEYWORDS = (
        "WAVE", "DOUBLE_WAVE", "ARC", "CHORD", "PIE",
        "ELLIPSE", "OVAL", "MOON", "TEAR", "TEARDROP", "CLOUD",
    )

    def is_arc_like(shape) -> bool:
        p = shape.prst.upper()
        return any(k in p for k in ARC_KEYWORDS)

    # 位置："从下至上页面底部 20% 高度范围内" → 形状整体位于页面纵向
    # 80%-100% 区间内。以形状顶部 s.y 落入区间为基准 —— 顶部在 80% 以下
    # 即整体在底部 20% 内。保留 2% 容差兼容办公软件 EMU→百分比的渲染精度
    # 损失。
    bottom_arcs = [
        s for s in slide
        if is_arc_like(s)
        and s.y / model.slide_h >= 0.80 - 0.02
    ]

    # 分别统计绿色和橙色弧形色块的填充面积（以外包围盒面积近似——
    # OOXML 层面无法访问曲线积分面积，各家办公软件亦以包围盒近似渲染）。
    green_shapes = [s for s in bottom_arcs if is_green(s.fill)]
    orange_shapes = [s for s in bottom_arcs if is_orange(s.fill)]
    green_area = sum(s.area for s in green_shapes)
    orange_area = sum(s.area for s in orange_shapes)

    # 有绿色弧形色块
    if green_area <= 0:
        return fail("第2页底部 20% 未找到绿色弧形色块")
    # 有橙色弧形色块
    if orange_area <= 0:
        return fail("第2页底部 20% 未找到橙色弧形色块")
    # 橙色面积小于绿色
    if orange_area >= green_area:
        return fail(
            f"第2页底部 20% 橙色弧形色块面积 {orange_area:.2f}cm² "
            f"未小于绿色 {green_area:.2f}cm²"
        )

    return ok(
        f"第2页底部 20% 检测到绿色/橙色弧形色块，"
        f"橙色面积 {orange_area:.2f}cm² 小于绿色 {green_area:.2f}cm²"
    )


def bottom_icon_with_label(model: PptModel, label: str, icon_chars: list[str]) -> CheckResult:
    word = first_text(model, 2, label)
    if not word:
        return fail(f"第2页底部未找到白色文本 {label}")
    if not (word.cy / model.slide_h >= 0.85 and any(is_white(c) for c in font_colors(word))):
        return fail(f"{label} 不在底部或不是白色")
    # 细则要求图标为"组合图形或者整体是一张图片"，且位于该文本"下方"对应的图标——
    # 即图标必须紧邻 word 上方的小邻域内（垂直差 ≤ 2cm，水平中心差 ≤ 0.8cm），
    # 避免把页面上半部的装饰元素错误计入。
    region_top = word.y - 2.0
    candidates = [s for s in model.slide(2)
                  if s is not word
                  and s.cy < word.cy and s.cy >= region_top
                  and abs(s.cx - word.cx) <= 0.8]
    # 路径 A：存在一张图片
    if any(is_picture(s) for s in candidates):
        return ok(f"第2页底部 {label} 上方找到对应图片")
    # 路径 B：存在 ≥2 个独立的"图形"（具备几何 prst 且有填充或线条），构成组合
    geom_parts = [s for s in candidates
                  if s.prst and s.prst not in ("RECTANGLE",)
                  and (s.fill is not None or s.line is not None)]
    if len(geom_parts) >= 2:
        return ok(f"第2页底部 {label} 上方找到对应组合图形（{len(geom_parts)} 件）")
    return fail(f"第2页底部 {label} 上方未找到符合细则的组合图形或图片")


def c33_sustainable(model: PptModel) -> CheckResult:
    """+3：第2页幻灯片左下角出现一个绿色的"带有白色脉络的叶子"组合图形
    或者整体是一个图片，其下方有白色文本"SUSTAINABLE"。
    """
    slide = model.slide(2)
    sw, sh = model.slide_w, model.slide_h

    # 1) 白色文本 "SUSTAINABLE"
    #    - 内容：包含 SUSTAINABLE
    #    - 颜色：白色（<a:solidFill><a:srgbClr val="FFFFFF"/>；PowerPoint /
    #      WPS / Keynote 均写入 srgbClr，`font_colors()` 归一化读取，跨办公
    #      软件有效）
    label = None
    for s in slide:
        if not s.text:
            continue
        if "SUSTAINABLE" not in s.text.upper():
            continue
        if not any(is_white(c) for c in font_colors(s)):
            continue
        label = s
        break
    if label is None:
        return fail('第2页未找到白色文本 "SUSTAINABLE"')

    # 2) 左下角：文本整体位于页面左半 (x2 ≤ 50%) 且底部 (y ≥ 60%)。
    #    ±2% 容差兼容 EMU→cm→百分比在各办公软件间的渲染精度差异。
    if not (label.x2 / sw <= 0.50 + 0.02 and label.y / sh >= 0.60 - 0.02):
        return fail(
            f'"SUSTAINABLE" 文本未位于第2页左下角（左至右 '
            f'{label.x/sw*100:.1f}%–{label.x2/sw*100:.1f}%，上至下 '
            f'{label.y/sh*100:.1f}%–{label.y2/sh*100:.1f}%）'
        )

    # 3) 其下方：图标位于 SUSTAINABLE 文本上方（icon 底部在 label 顶部
    #    之上或与文本中心线接近的紧邻位置），水平中心与文本接近，且同样
    #    位于左下角区域。相邻距离设为 2.5cm（用于容纳办公软件中 icon 与
    #    label 之间的行间距）。
    def is_upper_neighbor(s) -> bool:
        if s is label:
            return False
        if s.cy >= label.cy:  # 必须在文本上方
            return False
        if (label.cy - s.cy) > 2.5:  # 与文本相邻
            return False
        if abs(s.cx - label.cx) > 1.0:  # 与文本大致同列
            return False
        # 图标本身也应位于左下角区域
        if s.x / sw > 0.50 + 0.02:
            return False
        if s.y2 / sh < 0.60 - 0.02:
            return False
        return True

    upper = [s for s in slide if is_upper_neighbor(s)]
    if not upper:
        return fail('"SUSTAINABLE" 文本上方未找到图标形状')

    # 4) 图标判定（互斥两条路径，任一满足即通过）：
    #    路径 A —— 整体是一张图片：<p:pic>（在 python-pptx 中 shape.type
    #             为 PICTURE；WPS/Keynote 同样以 <p:pic> 表达）
    #    路径 B —— 绿色"带有白色脉络的叶子"组合图形：
    #             - ≥ 2 件独立几何形状（具备 prstGeom）
    #             - 至少一件绿色填充或轮廓（叶子主体）
    #             - 至少一件白色填充或轮廓（脉络）
    #    prstGeom 与 srgbClr 是 OOXML 标准，跨办公软件写法一致。

    # (A) 图片
    if any(is_picture(s) for s in upper):
        return ok('第2页左下角 "SUSTAINABLE" 上方找到图标图片')

    # (B) 绿色 + 白色脉络的组合图形
    parts = [
        s for s in upper
        if s.prst
        and (s.fill is not None or s.line is not None)
    ]
    has_green = any(is_green(s.fill) or is_green(s.line) for s in parts)
    has_white = any(is_white(s.fill) or is_white(s.line) for s in parts)
    if len(parts) >= 2 and has_green and has_white:
        return ok(
            f'第2页左下角 "SUSTAINABLE" 上方找到"绿色带白色脉络叶子"'
            f'组合图形（{len(parts)} 件）'
        )

    return fail(
        '第2页左下角 "SUSTAINABLE" 上方未找到"绿色带白色脉络叶子"'
        '组合图形或图片'
    )


def c34_innovative(model: PptModel) -> CheckResult:
    """+3：第2页幻灯片底部左侧出现一个由两个双向弯曲箭头组成的组合图形
    或者整体是一个图片，其下方有白色文本"INNOVATIVE"。
    """
    slide = model.slide(2)
    sw, sh = model.slide_w, model.slide_h

    # 1) 白色文本 "INNOVATIVE"
    label = None
    for s in slide:
        if not s.text:
            continue
        if "INNOVATIVE" not in s.text.upper():
            continue
        if not any(is_white(c) for c in font_colors(s)):
            continue
        label = s
        break
    if label is None:
        return fail('第2页未找到白色文本 "INNOVATIVE"')

    # 2) 底部：文本整体位于页面底部（y ≥ 60%）。水平位置 rubric 描述为
    #    "底部左侧"，但四个标签（SUSTAINABLE / INNOVATIVE / RELIABLE /
    #    PROFESSIONAL）实际横向平铺整个底部，因此这里只强制"底部"约束，
    #    不再强制"左侧"，避免误伤合规排版。±2% 容差兼容 EMU→cm→百分比
    #    在各办公软件间的渲染精度差异。
    if not (label.y / sh >= 0.60 - 0.02):
        return fail(
            f'"INNOVATIVE" 文本未位于第2页底部（上至下 '
            f'{label.y/sh*100:.1f}%–{label.y2/sh*100:.1f}%）'
        )

    # 3) 其下方：图标位于 INNOVATIVE 文本上方（icon 中心在 label 中心
    #    之上，垂直距离 ≤ 2.5cm，水平中心差 ≤ 1.0cm），且同样位于底部
    #    区域。此近邻窗口用于避免把页面上半部的装饰元素错误计入。
    def is_upper_neighbor(s: "ShapeInfo") -> bool:
        if s is label:
            return False
        if s.cy >= label.cy:  # 必须在文本上方
            return False
        if (label.cy - s.cy) > 2.5:  # 与文本相邻
            return False
        if abs(s.cx - label.cx) > 1.0:  # 与文本大致同列
            return False
        if s.y2 / sh < 0.60 - 0.02:  # 图标本身也应位于底部
            return False
        return True

    upper = [s for s in slide if is_upper_neighbor(s)]
    if not upper:
        return fail('"INNOVATIVE" 文本上方未找到图标形状')

    # 4) 图标判定（互斥两条路径，任一满足即通过）：
    #    路径 A —— 整体是一张图片：<p:pic>（python-pptx: PICTURE；WPS /
    #             Keynote 亦以 <p:pic> 表达）。图片本身是否为"双向弯曲
    #             箭头"无法在 OOXML 级别校验，rubric 允许"整体是一张
    #             图片"的写法，按图片存在即视为满足。
    #    路径 B —— 由两个双向弯曲箭头组成的组合图形：
    #             - "弯曲箭头"：OOXML prstGeom 明确的弯曲/回转/环形箭头
    #               预设，包括 curvedRightArrow / curvedLeftArrow /
    #               curvedUpArrow / curvedDownArrow / bentArrow /
    #               bentUpArrow / uturnArrow / circularArrow /
    #               leftCircularArrow / leftRightCircularArrow。以
    #               "ARROW" 关键字兜底并要求 CURVED/BENT/UTURN/
    #               CIRCULAR 之一，避免把 straight / left / right /
    #               leftRight 等直线箭头误判为弯曲。
    #             - "双向"：整体由两个弯曲箭头拼合，构成回环/相对指向；
    #               在 OOXML 中通常表达为两件独立的 curvedXxxArrow 形
    #               状，因此要求至少 2 件命中的箭头形状。
    #             - 若 OOXML 用 freeform 曲线绘制箭头（无 arrow prst），
    #               作为兜底：允许"箭头字符 ↻/↺/⇄/⇋/↔/arrow"文本 +
    #               ≥1 件 freeform 曲线 的组合命中。

    # (A) 图片
    if any(is_picture(s) for s in upper):
        return ok('第2页底部 "INNOVATIVE" 上方找到图标图片')

    # (B) 两个双向弯曲箭头组合
    CURVED_ARROW_KEYS = (
        "CURVEDRIGHTARROW", "CURVEDLEFTARROW",
        "CURVEDUPARROW", "CURVEDDOWNARROW",
        "BENTARROW", "BENTUPARROW", "UTURNARROW",
        "CIRCULARARROW", "LEFTCIRCULARARROW", "LEFTRIGHTCIRCULARARROW",
    )

    def _is_curved_arrow(s: "ShapeInfo") -> bool:
        if not s.prst:
            return False
        p = s.prst.upper()
        if "ARROW" not in p:
            return False
        # 精确匹配已知弯曲箭头预设
        if p in CURVED_ARROW_KEYS:
            return True
        # 关键字兜底：包含 CURVED / BENT / UTURN / CIRCULAR
        if any(k in p for k in ("CURVED", "BENT", "UTURN", "CIRCULAR")):
            return True
        return False

    curved_arrows: list[ShapeInfo] = [s for s in upper if _is_curved_arrow(s)]
    if len(curved_arrows) >= 2:
        return ok(
            f'第2页底部 "INNOVATIVE" 上方找到双向弯曲箭头组合（{len(curved_arrows)} 件）'
        )

    # 兜底：freeform 曲线绘制的箭头 —— 必须同时具备"箭头字符或 arrow 文本"
    # 语义标记（rubric 明确要求箭头，若纯 freeform 无任何箭头语义标记则视
    # 为不命中，避免任意曲线组合冒名顶替）。
    arrow_chars = ("↻", "↺", "⇄", "⇋", "↔", "⟳", "⟲")
    has_arrow_semantic = any(
        (s.text and (
            any(ch in s.text for ch in arrow_chars)
            or "ARROW" in s.text.upper()
        ))
        for s in upper
    )
    freeform_parts: list[ShapeInfo] = [
        s for s in upper
        if s.prst
        and ("FREEFORM" in s.prst.upper() or s.prst == "")
        and (s.fill is not None or s.line is not None)
    ]
    if has_arrow_semantic and len(freeform_parts) >= 1:
        return ok(
            f'第2页底部 "INNOVATIVE" 上方找到弯曲箭头 freeform 组合（{len(freeform_parts)} 件，含箭头语义文本）'
        )

    return fail(
        '第2页底部 "INNOVATIVE" 上方未找到"两个双向弯曲箭头"组合图形或图片'
    )


def c35_reliable(model: PptModel) -> CheckResult:
    """+3：第2页幻灯片底部左侧出现一个由盾牌和对勾组成的组合图形或者
    整体是一个图片，其下方有白色文本"RELIABLE"。
    """
    slide = model.slide(2)
    sw, sh = model.slide_w, model.slide_h

    # 1) 白色文本 "RELIABLE"
    #    内容含 RELIABLE，字体颜色为白色（OOXML <a:srgbClr val="FFFFFF"/>，
    #    PowerPoint / WPS / Keynote 均以此形式写入，`font_colors()` 已作跨
    #    办公软件的归一化读取）。
    label = None
    for s in slide:
        if not s.text:
            continue
        if "RELIABLE" not in s.text.upper():
            continue
        if not any(is_white(c) for c in font_colors(s)):
            continue
        label = s
        break
    if label is None:
        return fail('第2页未找到白色文本 "RELIABLE"')

    # 2) 底部左侧：文本整体位于页面底部 (y ≥ 60%) 且左侧 (x2 ≤ 50%)。
    #    ±2% 容差兼容 EMU→cm→百分比在各办公软件间的渲染精度差异。
    if not (label.y / sh >= 0.60 - 0.02 and label.x2 / sw <= 0.50 + 0.02):
        return fail(
            f'"RELIABLE" 文本未位于第2页底部左侧（左至右 '
            f'{label.x/sw*100:.1f}%–{label.x2/sw*100:.1f}%，上至下 '
            f'{label.y/sh*100:.1f}%–{label.y2/sh*100:.1f}%）'
        )

    # 3) "其下方有 RELIABLE" → 图标位于 RELIABLE 文本上方相邻位置，水平
    #    中心与文本对齐；图标本身也须位于底部左侧区域。
    def is_upper_neighbor(s) -> bool:
        if s is label:
            return False
        if s.cy >= label.cy:              # 图标必须在文本上方
            return False
        if (label.cy - s.cy) > 2.5:        # 与文本相邻
            return False
        if abs(s.cx - label.cx) > 1.0:     # 与文本大致同列
            return False
        if s.x / sw > 0.50 + 0.02:         # 图标本体亦在左侧
            return False
        if s.y2 / sh < 0.60 - 0.02:        # 图标本体亦在底部
            return False
        return True

    upper = [s for s in slide if is_upper_neighbor(s)]
    if not upper:
        return fail('"RELIABLE" 文本上方未找到图标形状')

    # 4) 图标判定（互斥两条路径，任一满足即通过）：
    #    路径 A —— 整体是一张图片：<p:pic>（PowerPoint / WPS / Keynote
    #             均以 <p:pic> 表达）。图片本身是否为"盾牌+对勾"无法在
    #             OOXML 级别校验，rubric 允许"整体是一张图片"的写法，按
    #             图片存在即视为满足。
    #    路径 B —— 由"盾牌"和"对勾"组成的组合图形：必须同时识别到
    #             「盾牌」与「对勾」两个语义构件，避免"文本上方有 ≥2 件
    #             任意几何"的宽松判定。
    #             - 盾牌：OOXML `prstGeom` 无原生 shield 预设，办公软件里
    #               通常采用 homePlate / pentagon / chevron / roundRect
    #               这类"上圆下尖"或"上方下尖"形状（PowerPoint 内置的
    #               "盾牌 SmartArt"/"盾牌"形状在导出为 OOXML 时会落到
    #               homePlate 或 freeform）。判定：
    #                 a) prst 含关键字 SHIELD；或
    #                 b) prst 属于 {homePlate, pentagon, chevron}；或
    #                 c) prst 为 freeform / 空 prst 且宽高比接近 1:1
    #                    (0.7 ≤ w/h ≤ 1.3) 且尺寸 ≥ 0.6cm（外轮廓兜底）。
    #             - 对勾：OOXML 无原生 checkmark 预设，办公软件通常用
    #               freeform 折线或 rightArrow 变体绘制。判定：
    #                 a) prst 含关键字 CHECK；或
    #                 b) 形状包含 "✓/✔/√/check" 文本字符；或
    #                 c) prst 为 freeform 且落入盾牌形状内部（面积 ≥ 50%
    #                    重叠或中心在盾牌包围盒内）作为兜底。

    # (A) 图片
    if any(is_picture(s) for s in upper):
        return ok('第2页底部左侧 "RELIABLE" 上方找到图标图片')

    # (B) 盾牌 + 对勾
    def _is_shield(s: "ShapeInfo") -> bool:
        if not s.prst:
            # 空 prst → freeform / custGeom 兜底：接近方形且尺寸够大
            if s.w > 0 and s.h > 0 and 0.7 <= s.w / s.h <= 1.3 \
                    and min(s.w, s.h) >= 0.6:
                return True
            return False
        p = s.prst.upper()
        if "SHIELD" in p:
            return True
        if p in ("HOMEPLATE", "PENTAGON", "CHEVRON"):
            return True
        if ("FREEFORM" in p) and s.w > 0 and s.h > 0 \
                and 0.7 <= s.w / s.h <= 1.3 and min(s.w, s.h) >= 0.6:
            return True
        return False

    def _is_check(s: "ShapeInfo", shield: "ShapeInfo | None" = None) -> bool:
        if s.prst and "CHECK" in s.prst.upper():
            return True
        if s.text and any(ch in s.text for ch in ("✓", "✔", "√")):
            return True
        if s.text and "CHECK" in s.text.upper():
            return True
        if shield is not None and s.prst and "FREEFORM" in s.prst.upper():
            if s.area > 0 and overlap_area(s, shield) / s.area >= 0.50:
                return True
            if (shield.x <= s.cx <= shield.x2
                    and shield.y <= s.cy <= shield.y2):
                return True
        return False

    shields: list[ShapeInfo] = [s for s in upper if _is_shield(s)]
    if shields:
        for sh in shields:
            checks: list[ShapeInfo] = [
                s for s in upper
                if s is not sh and _is_check(s, sh)
            ]
            if checks:
                return ok(
                    f'第2页底部左侧 "RELIABLE" 上方找到"盾牌 + 对勾"组合图形（盾牌 {len(shields)} 件、对勾 {len(checks)} 件）'
                )
        # 盾牌自身文本里带 ✓/✔/√/check（PowerPoint 允许盾牌形状承载对勾字符）
        for sh in shields:
            if sh.text and any(ch in sh.text for ch in ("✓", "✔", "√")):
                return ok(
                    '第2页底部左侧 "RELIABLE" 上方找到"盾牌 + 对勾"组合图形（盾牌形状内含对勾字符）'
                )

    return fail(
        '第2页底部左侧 "RELIABLE" 上方未找到"盾牌 + 对勾"组合图形或图片'
    )


def c36_global(model: PptModel) -> CheckResult:
    """+3：第2页幻灯片底部出现一个圆形和线条组成的组合图形或者整体是
    一个图片，其下方有白色文本"GLOBAL"。
    """
    slide = model.slide(2)
    sw, sh = model.slide_w, model.slide_h

    # 1) 白色文本 "GLOBAL"
    #    内容含 GLOBAL，字体颜色为白色（OOXML <a:srgbClr val="FFFFFF"/>，
    #    PowerPoint / WPS / Keynote 均以此形式写入，`font_colors()` 已作跨
    #    办公软件的归一化读取）。
    label = None
    for s in slide:
        if not s.text:
            continue
        if "GLOBAL" not in s.text.upper():
            continue
        if not any(is_white(c) for c in font_colors(s)):
            continue
        label = s
        break
    if label is None:
        return fail('第2页未找到白色文本 "GLOBAL"')

    # 2) 底部：文本整体位于页面底部 (y ≥ 60%)。
    #    ±2% 容差兼容 EMU→cm→百分比在各办公软件间的渲染精度差异。
    if not (label.y / sh >= 0.60 - 0.02):
        return fail(
            f'"GLOBAL" 文本未位于第2页底部（上至下 '
            f'{label.y/sh*100:.1f}%–{label.y2/sh*100:.1f}%）'
        )

    # 3) "其下方有 GLOBAL" → 图标位于 GLOBAL 文本上方相邻位置，水平中心
    #    与文本对齐；图标本身也须位于底部。
    def is_upper_neighbor(s) -> bool:
        if s is label:
            return False
        if s.cy >= label.cy:              # 图标必须在文本上方
            return False
        if (label.cy - s.cy) > 2.5:        # 与文本相邻
            return False
        if abs(s.cx - label.cx) > 1.0:     # 与文本大致同列
            return False
        if s.y2 / sh < 0.60 - 0.02:        # 图标本体亦在底部
            return False
        return True

    upper = [s for s in slide if is_upper_neighbor(s)]
    if not upper:
        return fail('"GLOBAL" 文本上方未找到图标形状')

    # 4) 图标判定（互斥两条路径，任一满足即通过）：
    #    路径 A —— 整体是一张图片：<p:pic>（PowerPoint / WPS / Keynote
    #             均以 <p:pic> 表达）。
    #    路径 B —— 由"圆形"和"线条"组成的组合图形：
    #             - 至少一件圆形（OOXML prst="ellipse"，python-pptx OVAL）
    #             - 至少一件线条（OOXML prst="line"，或宽/高之一为 0）
    #    prstGeom 是 OOXML 标准，跨办公软件写法一致。

    # (A) 图片
    if any(is_picture(s) for s in upper):
        return ok('第2页底部 "GLOBAL" 上方找到图标图片')

    # (B) 圆形 + 线条 组合图形
    has_oval = any(is_oval(s) for s in upper)
    has_line = any(is_line(s) for s in upper)
    if has_oval and has_line:
        return ok('第2页底部 "GLOBAL" 上方找到"圆形 + 线条"组合图形')

    return fail(
        '第2页底部 "GLOBAL" 上方未找到"圆形 + 线条"组合图形或图片'
    )


CRITERIA: list[Criterion] = [
    Criterion("D2-01", 5, "第1页左上角2-4cm logo组件/图片，含中上橙色和绿色组件", c01_logo_top_left),
    Criterion("D2-02", 1, "第1页中文公司名：微软雅黑/黑体24-32磅，顶部10%-25%", c02_cn_company),
    Criterion("D2-03", 1, "第2页英文公司名：Arial 24-32磅，两行，左0%-70%、顶10%-25%", c03_en_company_slide2),
    Criterion("D2-04", 1, "两页均有24-26cm×13-15cm圆角矩形边框，黑/灰/绿色实线", c04_borders),
    Criterion("D2-05", 1, "第1页中文公司名下方英文副标题：Arial 14-22磅橙色，位于顶部10%-25%", c05_en_company_slide1_subtitle),
    Criterion("D2-06", 3, "第1页英文副标题下方顶部23%-26%双橙线+叶片", c06_slide1_top_separator),
    Criterion("D2-07", 1, "两页右上角3-4cm二维码，宽高比1:1", c07_qr_both),
    Criterion("D2-08", 1, "第1页沈知行：左10%-45%、上30%-60%，宋/楷体加粗40-50磅绿色", c08_name_cn),
    Criterion("D2-09", 3, "第1页沈知行下方页面50%处双绿线+叶片", c09_name_separator),
    Criterion("D2-10", 1, "第1页M.P.：页面下方70%-80%、沈知行下方左对齐，Arial 13-20Pt绿色", c10_mp_slide1),
    Criterion("D2-11", 1, "两页电话文本位于M.P.右侧，Arial/Times 15-22Pt，黑/绿色", c11_phone_text_both),
    Criterion("D2-12", 3, "两页M.P.左侧深绿色0.8-1cm圆角矩形白色电话图案", c12_phone_icon_both),
    Criterion("D2-13", 1, "两页邮箱文本：页面70%-80%高、10%-40%宽，Arial 13-17磅", c13_email_text_both),
    Criterion("D2-14", 5, "两页邮箱左侧深绿色0.8-1cm圆角矩形白色信封图案", c14_email_icon_both),
    Criterion("D2-15", 1, "第1页中文地址：80%-90%高、10%-40%宽，邮箱下方左对齐，微软雅黑/黑体9-15磅", c15_cn_address_text),
    Criterion("D2-16", 3, "第1页中文地址左侧深绿色0.8-1cm圆角矩形白色定位图案", c16_cn_address_icon),
    Criterion("D2-17", 1, "两页45%-55%宽度处绿色竖向单实线，长6-8.5cm", c17_vertical_line_both),
    Criterion("D2-18", 3, "第1页绿色主营胶囊：20%-30%高、50%-75%宽，白色微软雅黑/黑体14-16磅", c18_cn_business_pill),
    Criterion("D2-19", 3, "第1页主营下方中文业务列表：逐行左对齐，标题绿13-15磅，条目黑11-13磅", c19_cn_business_list),
    Criterion("D2-20", 3, "第1页环保涂层左侧浅灰圆形及绿色水滴/叶子图案", c20_icon_eco_cn),
    Criterion("D2-21", 3, "第1页绿色材料左侧浅灰圆形及绿色叶子图案", c21_icon_green_cn),
    Criterion("D2-22", 3, "第1页应用方案左侧浅灰圆形（直径1.3-1.5cm）中间含绿色带白色箱盖/锁扣立体纸箱组合或图片", c22_icon_app_cn),
    Criterion("D2-23", 3, "第2页英文公司名左下方页面20%-50%双橙线（同一水平、单实线、长2.5-4cm）+ 橙/黄叶片组合或图片(0.3-0.4×0.4-0.5cm)", c23_slide2_title_separator),
    Criterion("D2-24", 1, "第2页Ethan Shen：左5%-45%、上30%-60%，Brush Script MT/Script MT Bold/Alex Brush/Segoe Script 25-35磅", c24_ethan),
    Criterion("D2-25", 1, "第2页邮箱下方英文地址三行：单独成行、左对齐，Arial 10-12磅，黑/绿/灰色", c25_en_address_lines),
    Criterion("D2-26", 3, "第2页英文地址末行左侧深绿色0.8-1cm圆角矩形白色定位图案", c26_en_address_icon),
    Criterion("D2-27", 3, "第2页英文公司名右下方绿色圆角矩形/胶囊：25%-40%高、50%-75%宽，内含MAIN BUSINESS 白色 Arial/黑体 10-15磅", c27_main_business_pill),
    Criterion("D2-28", 3, "第2页MAIN BUSINESS下方英文业务列表：单独成行、左对齐；标题绿Arial 11-13磅，条目黑/灰/绿 Arial 9-11磅", c28_en_business_list),
    Criterion("D2-29", 3, "第2页ECO COATINGS左侧浅灰圆形（直径1.3-1.5cm）中间含绿色水滴+小叶子组合或图片", c29_icon_eco_en),
    Criterion("D2-30", 3, "第2页GREEN MATERIALS左侧浅灰圆形（直径1.3-1.5cm）中间含绿色带白色脉络叶子组合或图片", c30_icon_green_en),
    Criterion("D2-31", 3, "第2页APPLICATION SOLUTIONS左侧浅灰圆形（直径1.3-1.5cm）中间含绿色带白色箱盖/锁扣立体纸箱组合或图片", c31_icon_app_en),
    Criterion("D2-32", 3, "第1页从下至上底部10%高度范围内有绿色和橙色的弧形色块，橙色面积小于绿色", c32_bottom_arcs_slide1),
    Criterion("D2-33", 3, "第2页左下角绿色带白色脉络叶子组合图形/图片，其下方白色文本SUSTAINABLE", c33_sustainable),
    Criterion("D2-34", 3, "第2页底部双向弯曲箭头/图片及白色INNOVATIVE", c34_innovative),
    Criterion("D2-35", 3, "第2页底部左侧盾牌+对勾组合图形/图片，其下方白色文本RELIABLE", c35_reliable),
    Criterion("D2-36", 3, "第2页底部圆形+线条组合图形/图片，其下方白色文本GLOBAL", c36_global),
    Criterion("D2-37", 3, "第2页从下至上底部20%高度范围内有绿色和橙色的弧形色块，橙色面积小于绿色", c37_bottom_arcs_slide2),
]


def _build_dim2_items(model: Optional[PptModel], dim1_pass: bool) -> tuple[list[dict[str, object]], int, int]:
    """按 CRITERIA 逐项检测，返回 (dim2_items, total_score, max_score)。

    - 命中项与未命中项均记录，便于外部按矩阵方式对齐。
    - 若维度一未通过 / model 为 None，则不进入维度二，返回空 dim2_items。
    """
    max_score = sum(c.score for c in CRITERIA)
    if not dim1_pass or model is None:
        return [], 0, max_score

    items: list[dict[str, object]] = []
    total = 0
    for criterion in CRITERIA:
        errored = False
        try:
            result = criterion.check(model)
        except Exception as exc:
            result = fail(f"检测异常：{exc}")
            errored = True
        hit = bool(result.passed)
        delta = criterion.score if hit else 0
        total += delta
        # errored 变量仍然保留，便于后续追踪评分异常；此处不再将其
        # 具体理由写入返回结构，`detail` 固定返回空字符串以满足对外
        # 接口"清空 detail"的要求。评分逻辑、命中判定与最终得分完全
        # 依赖 `hit` / `delta`，不受此调整影响。
        _ = errored
        items.append({
            "rule": criterion.title,
            "max_delta": criterion.score,
            "delta": delta,
            "hit": hit,
            "detail": "",
        })
    return items, total, max_score


def evaluate(dir_path: str) -> dict[str, object]:
    """统一评估入口。

    参数:
        dir_path: 脚本所在目录的路径（str）。脚本在该目录里定位并打开被评估的
            .pptx 文档（若存在多个则取修改时间最新的一个）。

    返回:
        参考《脚本接口差异与统一建议.md》§2.2 定义的结构化字典。
    """
    result: dict[str, object] = {
        "id": "043",
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": sum(c.score for c in CRITERIA),
    }

    try:
        directory = Path(dir_path)
        if not directory.is_dir():
            result["status"] = "error"
            result["error"] = f"目录不存在：{dir_path}"
            return result

        # 在目录中定位被评估的 PPTX 文档（只识别 .pptx，不再兼容老版 .ppt 二进制格式）
        candidates = sorted(
            [p for p in directory.iterdir()
             if p.is_file() and p.suffix.lower() == ".pptx"
             and not p.name.startswith("~$")],
            key=lambda p: p.stat().st_mtime,
            reverse=True,
        )
        if not candidates:
            result["status"] = "error"
            result["error"] = f"目录中未找到 .pptx 文档：{dir_path}"
            return result

        target = candidates[0]
        result["file_name"] = target.name

        # 维度一
        d1_passed, d1_details, model = check_dimension_1(target)
        result["dim1_pass"] = bool(d1_passed)
        if not d1_passed:
            # 汇总维度一未通过原因（取所有 ✗ 开头的行）
            reasons = [d for d in d1_details if d.startswith("✗")]
            result["dim1_reason"] = "；".join(reasons) if reasons else "维度一未通过"
            return result

        # 维度二
        items, total, max_score = _build_dim2_items(model, True)
        result["dim2_items"] = items
        result["total_score"] = total
        result["max_score"] = max_score
        return result

    except Exception as exc:
        result["status"] = "error"
        result["error"] = f"{type(exc).__name__}: {exc}"
        return result


if __name__ == "__main__":
    # 仅用于本地调试：默认传入脚本所在目录。
    debug_dir = sys.argv[1] if len(sys.argv) > 1 else str(Path(__file__).resolve().parent)
    print(json.dumps(evaluate(debug_dir), ensure_ascii=False, indent=2))
