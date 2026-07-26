# -*- coding: utf-8 -*-
"""
自动评估脚本：对 officeval_063 目录下的 PPT 文档进行打分。

评估逻辑：
  维度1（可用与可修改性）—— 任一硬性条件不满足，整份文件直接 0 分；
  维度2（完成度评分细则）—— 加分项需细则中“每一个点”都满足才累加；
                            扣分项只要细则中“任一点”命中即累减。

对外只暴露一个函数 `evaluate(dir_path: str) -> dict`：
  - dir_path 为脚本所在目录，脚本自行在其中定位被评估的 .pptx 文档；
  - 主结果不通过 print 输出，而是以结构化字典返回；
  - 顶层 try/except 兜底，脚本自身异常统一表达为 status="error"。
"""

import os
import sys
import json
import re
from lxml import etree

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

# --------------------------- 通用工具 ---------------------------

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "c": "http://schemas.openxmlformats.org/drawingml/2006/chart",
}


def emu_to_cm(emu):
    return (emu or 0) / 360000.0


def emu_to_inch(emu):
    return (emu or 0) / 914400.0


def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip("#")
    return tuple(int(hex_str[i:i + 2], 16) for i in (0, 2, 4))


def color_distance(c1_hex, c2_hex):
    """两个十六进制颜色的欧氏距离（0~441）。"""
    r1, g1, b1 = hex_to_rgb(c1_hex)
    r2, g2, b2 = hex_to_rgb(c2_hex)
    return ((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) ** 0.5


def is_blueish(hex_color):
    """判断是否为浅蓝/青色系（B 通道相对较强、整体偏亮）。"""
    r, g, b = hex_to_rgb(hex_color)
    brightness = (r + g + b) / 3
    return brightness > 200 and b >= r and (b + g) / 2 > r


def is_beige(hex_color):
    """判断是否为米黄色系（R≈G > B，且偏暖）。"""
    r, g, b = hex_to_rgb(hex_color)
    return r > 200 and g > 190 and b < r and b < g and (r - b) > 15


# --------------------------- 维度1 检查 ---------------------------


def check_dim1(pptx_path: str):
    """维度1：交付文件为 .pptx 格式，能够正常打开。

    只做两件事：
      1) 后缀是 .pptx；
      2) 能够被 python-pptx 正常打开。
    其它（页数、可编辑、放映等）不再作为维度一的判据。
    """
    fails = []

    # 1. 后缀
    if not pptx_path.lower().endswith(".pptx"):
        fails.append("文件格式不是 .pptx")
        return False, fails

    # 2. 能否打开
    try:
        _ = Presentation(pptx_path)
    except Exception as e:
        fails.append(f"无法打开 PPT：{e}")
        return False, fails

    return True, []


# --------------------------- 维度2 工具函数 ---------------------------


def get_chart_on_slide(slide):
    for shp in slide.shapes:
        if getattr(shp, "has_chart", False):
            return shp.chart
    return None


def iter_runs(shape):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            yield p, r


def get_run_color_hex(run):
    """尽量取出 run 的 RGB 十六进制（大写）；取不到返回 None。"""
    try:
        c = run.font.color
        if c and c.rgb is not None:
            return str(c.rgb).upper()
    except Exception:
        pass
    # 退化：从 XML 直接抓 solidFill
    rPr = run._r.find(".//a:rPr", NS)
    if rPr is not None:
        clr = rPr.find(".//a:solidFill/a:srgbClr", NS)
        if clr is not None:
            return clr.get("val", "").upper()
    return None


def slide_bg_hex(slide):
    """兼容旧调用点：仅返回 slide 级 p:bg 内的 a:srgbClr 十六进制值（大写）。

    新代码请使用 :func:`resolve_slide_background_hex`，它会补充解析主题色/
    layout/master 背景、bgRef、gradFill、以及全页背景形状。
    """
    bg = slide._element.find(".//p:cSld/p:bg", NS)
    if bg is None:
        return None
    clr = bg.find(".//a:srgbClr", NS)
    return clr.get("val").upper() if clr is not None else None


# ---------- 背景色的完整 OOXML 解析（用于 +3 第5页背景规则） ---------- #

_PRESET_COLOR_HEX = {
    # 常用命名色 → 十六进制（不完全，够用即可，覆盖 rubric 场景）
    "beige": "F5F5DC",
    "bisque": "FFE4C4",
    "cornsilk": "FFF8DC",
    "wheat": "F5DEB3",
    "lightblue": "ADD8E6",
    "lightcyan": "E0FFFF",
    "aliceblue": "F0F8FF",
    "azure": "F0FFFF",
    "powderblue": "B0E0E6",
    "skyblue": "87CEEB",
    "lightsteelblue": "B0C4DE",
    "white": "FFFFFF",
    "black": "000000",
}


def _apply_lum_shade(rgb, lum_mod, lum_off, shade, tint):
    """按 OOXML 的 lumMod/lumOff/shade/tint（千分制值，0..100000）近似还原 RGB。"""
    r, g, b = rgb
    # lumMod: 亮度乘以百分比
    if lum_mod is not None:
        factor = lum_mod / 100000.0
        r = r * factor
        g = g * factor
        b = b * factor
    # lumOff: 亮度加偏移（相对 255 * off/100000）
    if lum_off is not None:
        off = lum_off / 100000.0 * 255.0
        r += off
        g += off
        b += off
    # shade: 变暗
    if shade is not None:
        factor = shade / 100000.0
        r *= factor
        g *= factor
        b *= factor
    # tint: 与白色混合
    if tint is not None:
        factor = tint / 100000.0
        r = r + (255 - r) * factor
        g = g + (255 - g) * factor
        b = b + (255 - b) * factor
    r = int(max(0, min(255, round(r))))
    g = int(max(0, min(255, round(g))))
    b = int(max(0, min(255, round(b))))
    return (r, g, b)


def _theme_scheme_map(prs):
    """从主题包(theme1.xml)中读取 a:clrScheme → {schemeClr键: hex}。

    schemeClr 的常见 val：bg1/bg2/tx1/tx2/accent1..6/hlink/folHlink。
    theme 中对应的元素名（顺序）：lt1/dk1/lt2/dk2/accent1..6/hlink/folHlink，
    以及 bg1<->lt1, tx1<->dk1, bg2<->lt2, tx2<->dk2 的等价映射。
    """
    scheme = {}
    theme_root = None
    for master in prs.slide_masters:
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                try:
                    theme_root = etree.fromstring(rel.target_part.blob)
                    break
                except Exception:
                    continue
        if theme_root is not None:
            break
    if theme_root is None:
        return scheme

    clr_scheme = theme_root.find(".//a:clrScheme", NS)
    if clr_scheme is None:
        return scheme

    # 对每一个 clr 元素解析其 srgbClr / sysClr / prstClr
    for child in clr_scheme:
        tag = etree.QName(child.tag).localname
        srgb = child.find("a:srgbClr", NS)
        sysc = child.find("a:sysClr", NS)
        prst = child.find("a:prstClr", NS)
        hex_val = None
        if srgb is not None:
            hex_val = (srgb.get("val") or "").upper()
        elif sysc is not None:
            hex_val = (sysc.get("lastClr") or "").upper()
        elif prst is not None:
            hex_val = _PRESET_COLOR_HEX.get((prst.get("val") or "").lower())
        if hex_val:
            scheme[tag] = hex_val
    # 等价别名
    aliases = {"bg1": "lt1", "bg2": "lt2", "tx1": "dk1", "tx2": "dk2"}
    for a, b in aliases.items():
        if a not in scheme and b in scheme:
            scheme[a] = scheme[b]
        if b not in scheme and a in scheme:
            scheme[b] = scheme[a]
    return scheme


def _resolve_color_element(color_el, scheme_map):
    """把一个 color 容器元素 (a:solidFill 或 a:srgbClr/a:schemeClr/a:prstClr 的父节点)
    解析成 RGB 元组，处理 lumMod/lumOff/shade/tint 修饰。找不到返回 None。
    """
    if color_el is None:
        return None
    # 定位实际颜色节点
    srgb = color_el.find(".//a:srgbClr", NS)
    scheme = color_el.find(".//a:schemeClr", NS)
    sysc = color_el.find(".//a:sysClr", NS)
    prst = color_el.find(".//a:prstClr", NS)
    chosen = srgb if srgb is not None else scheme if scheme is not None else sysc if sysc is not None else prst
    if chosen is None:
        return None

    def _get_int(el, name):
        node = el.find(f"a:{name}", NS)
        if node is not None and node.get("val") is not None:
            try:
                return int(node.get("val"))
            except Exception:
                return None
        return None

    lum_mod = _get_int(chosen, "lumMod")
    lum_off = _get_int(chosen, "lumOff")
    shade = _get_int(chosen, "shade")
    tint = _get_int(chosen, "tint")

    if chosen is srgb:
        try:
            hx = (chosen.get("val") or "").upper()
            base = hex_to_rgb(hx)
        except Exception:
            return None
    elif chosen is sysc:
        hx = (chosen.get("lastClr") or "").upper()
        if not hx:
            return None
        base = hex_to_rgb(hx)
    elif chosen is prst:
        hx = _PRESET_COLOR_HEX.get((chosen.get("val") or "").lower())
        if not hx:
            return None
        base = hex_to_rgb(hx)
    else:  # schemeClr
        key = chosen.get("val") or ""
        hx = scheme_map.get(key)
        if not hx:
            return None
        base = hex_to_rgb(hx)

    return _apply_lum_shade(base, lum_mod, lum_off, shade, tint)


def _extract_bg_rgb_from_bg_element(bg_el, scheme_map):
    """从 p:bg 元素抽出背景 RGB。支持 p:bgPr/a:solidFill、p:bgRef、a:gradFill。返回 (r,g,b) 或 None。"""
    if bg_el is None:
        return None
    # p:bgPr → a:solidFill
    solid = bg_el.find(".//p:bgPr/a:solidFill", NS)
    if solid is not None:
        return _resolve_color_element(solid, scheme_map)
    # p:bgPr → a:gradFill：取第一/最后色标里最亮者作为主色
    grad = bg_el.find(".//p:bgPr/a:gradFill", NS)
    if grad is not None:
        stops = grad.findall(".//a:gsLst/a:gs", NS)
        rgbs = []
        for gs in stops:
            rgb = _resolve_color_element(gs, scheme_map)
            if rgb:
                rgbs.append(rgb)
        if rgbs:
            # 选择"亮度中位"色作为主色，避免过深/过浅极端
            rgbs.sort(key=lambda t: sum(t))
            return rgbs[len(rgbs) // 2]
    # p:bgRef → 引用主题色
    bgref = bg_el.find(".//p:bgRef", NS)
    if bgref is not None:
        return _resolve_color_element(bgref, scheme_map)
    # p:bgPr → a:blipFill：图片背景，无稳定主色 → 交由上层用形状扫描兜底
    return None


def _fullpage_shape_bg_rgb(slide, scheme_map):
    """扫描 slide 上覆盖近整页的形状；若存在实心填充，返回该主色 RGB。"""
    sw = slide.part.package.presentation_part.presentation.slide_width
    sh = slide.part.package.presentation_part.presentation.slide_height
    page_area = max(sw * sh, 1)
    for shp in slide.shapes:
        w = shp.width or 0
        h = shp.height or 0
        if w * h < page_area * 0.9:
            continue
        # 需为实心填充：a:solidFill；避免文本框/带内容的形状
        sp_pr = None
        try:
            sp_pr = shp._element.find(".//p:spPr", NS)
        except Exception:
            sp_pr = None
        if sp_pr is None:
            # 某些形状使用 wp:spPr / 直接 a:spPr —— 兜底
            sp_pr = shp._element.find(".//a:spPr", NS)
        if sp_pr is None:
            continue
        solid = sp_pr.find("a:solidFill", NS)
        if solid is None:
            continue
        rgb = _resolve_color_element(solid, scheme_map)
        if rgb is not None:
            return rgb
    return None


def resolve_slide_background_hex(prs, slide):
    """返回该幻灯片"生效背景色"的十六进制 (RRGGBB 大写)，找不到返回 None。

    解析顺序（模仿 PowerPoint 渲染背景的选取过程）：
      1) slide 级 p:cSld/p:bg
      2) slide 级覆盖近整页的实色形状（视觉背景）
      3) slide layout 级 p:cSld/p:bg
      4) slide master 级 p:cSld/p:bg
    颜色源支持 srgbClr/schemeClr/sysClr/prstClr + lumMod/lumOff/shade/tint。
    """
    scheme_map = _theme_scheme_map(prs)

    # 1) slide bg
    bg_el = slide._element.find(".//p:cSld/p:bg", NS)
    rgb = _extract_bg_rgb_from_bg_element(bg_el, scheme_map)
    if rgb is not None:
        return "%02X%02X%02X" % rgb

    # 2) 全页背景形状（兜底"以形状充当背景"）
    rgb = _fullpage_shape_bg_rgb(slide, scheme_map)
    if rgb is not None:
        return "%02X%02X%02X" % rgb

    # 3) layout bg
    layout = slide.slide_layout
    bg_el = layout.element.find(".//p:cSld/p:bg", NS)
    rgb = _extract_bg_rgb_from_bg_element(bg_el, scheme_map)
    if rgb is not None:
        return "%02X%02X%02X" % rgb

    # 4) master bg
    master = layout.slide_master
    bg_el = master.element.find(".//p:cSld/p:bg", NS)
    rgb = _extract_bg_rgb_from_bg_element(bg_el, scheme_map)
    if rgb is not None:
        return "%02X%02X%02X" % rgb

    return None


# --------------------------- 维度2 逐条规则 ---------------------------


def rule_10_chart(prs):
    """+5  第10页曲线图对象 —— 严格按打分细则逐点校验：
      (1) 在结束页右侧或中部空白区域
      (2) 可编辑的折线图或平滑曲线图（原生 chart）
      (3) 含 "营收数值" 和 "盈利比值" 两个数据系列
      (4) 两条曲线颜色不同
      (5) 横轴从左至右依次为 "1月、2月、3月、4月、5月"
      (6) 营收数值五个月数据依次 68000、352000、486000、423000、579000
      (7) 盈利比值五个月数据依次 1.42、1.56、1.71、1.63、1.85
      (8) 准确显示 "营收数值"、"盈利比值"（图例文字）
      (9) 图例颜色与两条曲线一一对应
    任何一点不满足即不加分。

    严格化（回应反馈）：
      · (1) 位置判定不再只看图表中心 x 是否落在右/中部；还必须验证图表矩形
            与结束页上其它非背景内容元素**不相交**（占用面积 ≤ 图表面积 5%），
            以确认所处区域为"空白区域"。
      · (9) 图例↔曲线颜色一一对应不再只看 legendEntry 覆写颜色的唯一性和
            c:idx 唯一性；改为解析每个系列(c:order/c:idx)与其显示名称及
            线条颜色的实际映射，并检查图例条目未被 c:delete 删除。这样，
            即便某系列被隐藏（图例项被删除），或图例文字通过 legendEntry
            重命名/串位，也能被识别出来。
    """
    slide10 = prs.slides[9]
    sw = prs.slide_width
    sh = prs.slide_height

    # ---- (2) 原生可编辑图表对象 ----
    chart_shape = None
    for shp in slide10.shapes:
        if getattr(shp, "has_chart", False):
            chart_shape = shp
            break
    if chart_shape is None:
        return False, "未检测到第 10 页的原生图表对象（非可编辑图表）"

    chart = chart_shape.chart
    ct_name = str(chart.chart_type).upper()
    if "LINE" not in ct_name:
        return False, f"图表类型不是折线/平滑曲线图（实际 {chart.chart_type}）"

    # ---- (1) 位置：右侧或中部 + 所在矩形为空白区域 ----
    chart_L = chart_shape.left or 0
    chart_T = chart_shape.top or 0
    chart_W = chart_shape.width or 0
    chart_H = chart_shape.height or 0
    chart_R = chart_L + chart_W
    chart_B = chart_T + chart_H
    cx = chart_L + chart_W / 2
    cy = chart_T + chart_H / 2

    # 必须在页面内
    if not (0 <= cx <= sw and 0 <= cy <= sh):
        return False, "图表不在结束页页面范围内"
    in_right = cx >= sw * 0.5
    in_center = (sw * 1 / 3) <= cx <= (sw * 2 / 3)
    if not (in_right or in_center):
        return False, f"图表中心 x={emu_to_cm(cx):.1f}cm 既不在页面右侧也不在中部"

    # "空白区域" —— 图表矩形与结束页任何"非背景内容元素"的相交面积必须显著为 0。
    #   忽略：图表自身；覆盖近整页的装饰底图/背景形状（面积 ≥ 页面 95%）；
    #         无文本无可见几何轮廓的空形状（宽或高为 0）。
    chart_area = max(chart_W * chart_H, 1)
    page_area = max(sw * sh, 1)
    overlap_details = []
    for shp in slide10.shapes:
        if shp is chart_shape:
            continue
        sl = shp.left or 0
        st = shp.top or 0
        sw_ = shp.width or 0
        sh_ = shp.height or 0
        if sw_ <= 0 or sh_ <= 0:
            continue
        # 忽略覆盖近整页的装饰底图/背景
        if (sw_ * sh_) >= page_area * 0.95:
            continue
        # 空白文本框（无文本内容）不视作占位内容
        if shp.has_text_frame and not shp.text_frame.text.strip():
            continue
        sr_ = sl + sw_
        sb_ = st + sh_
        ix = max(0, min(chart_R, sr_) - max(chart_L, sl))
        iy = max(0, min(chart_B, sb_) - max(chart_T, st))
        inter = ix * iy
        if inter <= 0:
            continue
        # 相交面积超过图表面积 5% 视为占用图表所在区域
        if inter / chart_area > 0.05:
            overlap_details.append(
                f"{shp.name}(相交 {inter / chart_area * 100:.0f}%图表面积)"
            )
    if overlap_details:
        return False, "图表所在区域并非空白，与内容形状重叠：" + "; ".join(overlap_details[:3])

    # ---- (5) 横轴类别顺序 ----
    try:
        categories = list(chart.plots[0].categories)
    except Exception as e:
        return False, f"无法读取图表类别（横轴）：{e}"
    expected_cats = ["1月", "2月", "3月", "4月", "5月"]
    cat_strs = [str(c).strip() for c in categories]
    if cat_strs != expected_cats:
        return False, f"横轴月份/顺序不正确，期望 {expected_cats}，实际 {cat_strs}"

    # ---- (3) 两个数据系列：名称必须为 营收数值 与 盈利比值 ----
    series_list = list(chart.series)
    if len(series_list) != 2:
        return False, f"系列数量不是 2（实际 {len(series_list)}）"
    series_names = [(s.name or "").strip() for s in series_list]
    if set(series_names) != {"营收数值", "盈利比值"}:
        return False, f"系列名称不正确，需要 营收数值 与 盈利比值，实际 {series_names}"

    # ---- (6)(7) 系列数据精确比对 ----
    exp_rev = [68000, 352000, 486000, 423000, 579000]
    exp_ratio = [1.42, 1.56, 1.71, 1.63, 1.85]
    series_by_name = {(s.name or "").strip(): list(s.values) for s in series_list}
    rev = series_by_name["营收数值"]
    ratio = series_by_name["盈利比值"]
    if len(rev) != 5 or any(abs(a - b) > 0.5 for a, b in zip(rev, exp_rev)):
        return False, f"营收数值数据不正确，期望 {exp_rev}，实际 {rev}"
    if len(ratio) != 5 or any(abs(a - b) > 0.01 for a, b in zip(ratio, exp_ratio)):
        return False, f"盈利比值数据不正确，期望 {exp_ratio}，实际 {ratio}"

    # ---- (4) 两条曲线颜色不同 ----
    chart_xml = chart._chartSpace
    ser_nodes = chart_xml.findall(".//c:ser", NS)
    if len(ser_nodes) != 2:
        return False, f"chart XML 中 c:ser 数量异常：{len(ser_nodes)}"

    def _series_line_color(ser):
        """提取 series 的折线主颜色（按 OOXML 优先级）。返回 ('srgb'|'scheme'|'auto', value) 或 None。"""
        # 顺序：c:ser/c:spPr/a:ln 下 solidFill；其次 c:spPr/a:solidFill
        ln = ser.find(".//c:spPr/a:ln", NS)
        if ln is not None:
            srgb = ln.find(".//a:solidFill/a:srgbClr", NS)
            if srgb is not None:
                return ("srgb", srgb.get("val", "").upper())
            sch = ln.find(".//a:solidFill/a:schemeClr", NS)
            if sch is not None:
                return ("scheme", sch.get("val", ""))
        sp = ser.find(".//c:spPr", NS)
        if sp is not None:
            srgb = sp.find(".//a:solidFill/a:srgbClr", NS)
            if srgb is not None:
                return ("srgb", srgb.get("val", "").upper())
            sch = sp.find(".//a:solidFill/a:schemeClr", NS)
            if sch is not None:
                return ("scheme", sch.get("val", ""))
        # 找不到显式颜色（依赖默认主题）：以 series index 区分（PPT 默认每个 series 颜色不同）
        idx = ser.find("c:idx", NS)
        return ("auto", idx.get("val") if idx is not None else None)

    def _series_meta(ser):
        idx_el = ser.find("c:idx", NS)
        order_el = ser.find("c:order", NS)
        name_el = ser.find(".//c:tx//c:v", NS)
        return {
            "idx": int(idx_el.get("val")) if idx_el is not None and idx_el.get("val") else None,
            "order": int(order_el.get("val")) if order_el is not None and order_el.get("val") else None,
            "name": (name_el.text.strip() if name_el is not None and name_el.text else ""),
            "color": _series_line_color(ser),
        }

    series_meta = [_series_meta(s) for s in ser_nodes]
    line_colors = [m["color"] for m in series_meta]
    if any(c is None for c in line_colors):
        return False, "无法解析两条曲线的颜色"
    if line_colors[0] == line_colors[1]:
        return False, f"两条曲线颜色相同：{line_colors[0]}"

    # ---- 图例存在 ----
    if not chart.has_legend:
        return False, "图表未启用图例"

    # ---- (8) 图例文字准确显示 "营收数值"、"盈利比值" ----
    legend_texts = []
    for t in chart_xml.findall(".//c:legend//a:t", NS):
        if t.text:
            legend_texts.append(t.text.strip())
    legend_pool = set(legend_texts) | {m["name"] for m in series_meta}
    if "营收数值" not in legend_pool or "盈利比值" not in legend_pool:
        return False, (
            f"图例文本未准确显示 营收数值/盈利比值（图例文字={legend_texts}, "
            f"系列名={[m['name'] for m in series_meta]}）"
        )

    # ---- (9) 图例颜色与两条曲线一一对应 ----
    # 在 OOXML 中，legend 默认按 c:order 依次渲染每个 series 一个条目，
    # 标记色 = series 线条色。可能破坏"一一对应"的实际方式只有两类：
    #   A. c:legendEntry 用 c:delete val="1" 删除某条条目 → 只剩一条，无法一一对应
    #   B. c:legendEntry 通过 c:txPr 覆写文本颜色使两条条目看起来变成同色
    # 若两 series 的 order/idx 冲突（重复），说明源数据本身就已经不能与图例
    # 一一对应；一并作为失败条件。
    orders = [m["order"] for m in series_meta]
    idxs = [m["idx"] for m in series_meta]
    if len(set(orders)) != len(orders) or len(set(idxs)) != len(idxs):
        return False, f"系列 c:idx/c:order 存在重复，图例无法与曲线一一对应：idxs={idxs}, orders={orders}"

    deleted_entry_idxs = []
    override_text_colors: "dict[int, str]" = {}
    for le in chart_xml.findall(".//c:legend/c:legendEntry", NS):
        idx_el = le.find("c:idx", NS)
        entry_idx = int(idx_el.get("val")) if (idx_el is not None and idx_el.get("val")) else None
        delete = le.find("c:delete", NS)
        if delete is not None and delete.get("val", "1") in ("1", "true"):
            deleted_entry_idxs.append(entry_idx)
        # 覆写文本颜色（可能把两条图例文字都染成同色，视觉上无法区分）
        srgb = le.find(".//c:txPr//a:solidFill/a:srgbClr", NS)
        if srgb is not None and entry_idx is not None:
            override_text_colors[entry_idx] = srgb.get("val", "").upper()
    if deleted_entry_idxs:
        return False, f"图例条目被删除：idx={deleted_entry_idxs}，与曲线无法一一对应"
    if override_text_colors and len(set(override_text_colors.values())) < len(override_text_colors):
        return False, f"图例条目文本颜色被覆写为相同值：{override_text_colors}"

    # 构造对应关系摘要（系列名 → 线条色）便于回看
    mapping = ", ".join(f"{m['name']}→{m['color']}" for m in series_meta)

    return True, (
        f"位置={'右侧' if in_right else '中部'}空白区; 类型={chart.chart_type}; "
        f"系列={series_names}; 横轴={cat_strs}; "
        f"营收={rev}; 盈利={ratio}; 双色映射={{{mapping}}}; "
        f"图例文字含 营收数值/盈利比值，且颜色/条目与曲线一一对应"
    )


def rule_2_paiban(prs):
    """+1  第2页右侧第3个岗位卡片：空白处有标题 "排班岗"，
         字体为深绿色微软雅黑 24磅 加粗。逐点校验：
      (1) 位于第 2 页右侧第 3 个岗位卡片（按卡片中心 x 升序排列后取第 3 张，且其位于右侧）
      (2) 该卡片"空白处"有标题文本 "排班岗"
      (3) 字体 = 微软雅黑（Microsoft YaHei / 微软雅黑）
      (4) 颜色 = 深绿色
      (5) 字号 = 24 磅
      (6) 加粗
    任一不满足即不加分。
    """
    s2 = prs.slides[1]
    sw = prs.slide_width

    # ---- (1) 找出 3 张岗位卡片（高度较大的矩形容器），按中心 x 排序，取第 3 张 ----
    # 岗位卡片在本 PPT 中为高度 ≥ 3 inch 的 AUTO_SHAPE 容器
    # 注意：本页内容框（如 Text 5/Text 8）同样是 AUTO_SHAPE 且尺寸接近卡片外框，
    # 需排除"被另一个更大形状几乎完全包住"的内嵌内容框，避免索引错位。
    candidates = []
    for shp in s2.shapes:
        if shp.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        w = shp.width or 0
        h = shp.height or 0
        # 卡片特征：宽 ≥ 2 inch、高 ≥ 3 inch
        if w >= Emu(2 * 914400) and h >= Emu(3 * 914400):
            candidates.append(shp)

    def _rect(shp):
        l = shp.left or 0
        t = shp.top or 0
        return l, t, l + (shp.width or 0), t + (shp.height or 0)

    def _contained_in(inner, outer):
        # inner 的矩形几乎完全位于 outer 内部（允许 1mm 误差）
        if inner is outer:
            return False
        tol = Emu(int(0.1 * 360000))  # 0.1cm
        il, it, ir, ib = _rect(inner)
        ol, ot, or_, ob = _rect(outer)
        return (il >= ol - tol and it >= ot - tol
                and ir <= or_ + tol and ib <= ob + tol)

    cards = []
    for shp in candidates:
        # 若该候选形状被其他更大候选完全包住，则视为内容框，跳过
        nested = any(
            _contained_in(shp, other)
            and (other.width or 0) * (other.height or 0)
                > (shp.width or 0) * (shp.height or 0)
            for other in candidates
        )
        if not nested:
            cards.append(shp)
    cards.sort(key=lambda s: (s.left or 0) + (s.width or 0) / 2)

    if len(cards) < 3:
        return False, f"第 2 页未识别到 3 张岗位卡片（识别到 {len(cards)} 张）"

    third_card = cards[2]
    # 第 3 张必须在页面右侧
    third_cx = (third_card.left or 0) + (third_card.width or 0) / 2
    if third_cx < sw * 0.5:
        return False, f"第 3 张岗位卡片中心 x={emu_to_cm(third_cx):.1f}cm，不在右侧"

    # ---- (2) 在第 3 张卡片范围内查找文本 "排班岗"，且必须位于该卡片"空白处" ----
    tcx0 = third_card.left or 0
    tcy0 = third_card.top or 0
    tcx1 = tcx0 + (third_card.width or 0)
    tcy1 = tcy0 + (third_card.height or 0)

    # "空白处" 判定 —— 收集第 3 张卡片范围内、除标题候选之外的既有内容：
    #   · 卡片内的其它文本框（非空文本）
    #   · 图片 / 图标 / 自选图形（AUTO_SHAPE/FREEFORM/PICTURE/GROUP 等有实际视觉占位的形状）
    # 卡片外框（third_card 本身）不计入内容占位；标题候选自身自然也不算重叠。
    def _rect_of(shp):
        l = shp.left or 0
        t = shp.top or 0
        return l, t, l + (shp.width or 0), t + (shp.height or 0)

    def _center_inside_card(shp):
        scx = (shp.left or 0) + (shp.width or 0) / 2
        scy = (shp.top or 0) + (shp.height or 0) / 2
        return tcx0 <= scx <= tcx1 and tcy0 <= scy <= tcy1

    def _rects_overlap(a, b, tol=0):
        # a/b: (l, t, r, b)  tol 允许的最大间距（EMU）
        return not (a[2] + tol <= b[0] or b[2] + tol <= a[0]
                    or a[3] + tol <= b[1] or b[3] + tol <= a[1])

    existing_content_rects = []  # (rect, tag) 已有非空文本 / 视觉占位形状
    for shp in s2.shapes:
        if shp is third_card:
            continue
        if not _center_inside_card(shp):
            continue
        # 文本内容
        if shp.has_text_frame:
            txt = shp.text_frame.text.strip()
            if txt and txt != "排班岗":
                existing_content_rects.append((_rect_of(shp), f"卡内文本'{txt[:6]}'"))
                continue
            if not txt:
                # 空文本框不视作占位内容
                continue
        # 视觉占位（图片/图标/自选图形/成组）
        if shp.shape_type in (
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.FREEFORM,
            MSO_SHAPE_TYPE.GROUP,
        ):
            # 排除卡片自身（前面已 continue）；把明显是卡片主体轮廓的形状排除
            w = shp.width or 0
            h = shp.height or 0
            card_w = (third_card.width or 1)
            card_h = (third_card.height or 1)
            if w >= card_w * 0.85 and h >= card_h * 0.85:
                # 与卡片外框几乎重合的形状，视作卡片主体，不计入内容占位
                continue
            existing_content_rects.append((_rect_of(shp), f"卡内形状'{shp.name}'"))

    target_shp = None
    target_run = None
    blocked_reason = ""
    for shp in s2.shapes:
        if shp is third_card or not shp.has_text_frame:
            continue
        if not _center_inside_card(shp):
            continue
        if shp.text_frame.text.strip() != "排班岗":
            continue
        # ---- "空白处" 检查：标题矩形与卡片内已有内容不得显著重叠 ----
        title_rect = _rect_of(shp)
        title_area = max(
            (title_rect[2] - title_rect[0]) * (title_rect[3] - title_rect[1]), 1
        )
        overlap_hits = []
        for rc, tag in existing_content_rects:
            if not _rects_overlap(title_rect, rc):
                continue
            ix = max(0, min(title_rect[2], rc[2]) - max(title_rect[0], rc[0]))
            iy = max(0, min(title_rect[3], rc[3]) - max(title_rect[1], rc[1]))
            inter = ix * iy
            # 相交面积 > 标题面积 10% 即视为覆盖已有内容（非空白处）
            if inter / title_area > 0.10:
                overlap_hits.append(
                    f"{tag}(相交 {inter / title_area * 100:.0f}%标题面积)"
                )
        if overlap_hits:
            blocked_reason = (
                "标题'排班岗'未位于卡片空白处，与已有内容重叠：" + "; ".join(overlap_hits[:3])
            )
            continue

        target_shp = shp
        for _, r in iter_runs(shp):
            if r.text.strip() == "排班岗":
                target_run = r
                break
        if target_run is not None:
            break

    if target_run is None:
        if blocked_reason:
            return False, blocked_reason
        return False, "第 2 页右侧第 3 个岗位卡片内未找到标题 '排班岗'"

    # ---- (3)(4)(5)(6) 字体属性 ----
    font = target_run.font
    name = font.name or ""
    size_pt = font.size.pt if font.size else None
    bold = bool(font.bold)
    color_hex = get_run_color_hex(target_run) or ""

    reasons = []
    # (3) 微软雅黑
    if not ("YaHei" in name or "Microsoft YaHei" in name or "微软雅黑" in name):
        reasons.append(f"字体非微软雅黑（实际 '{name}'）")
    # (5) 24 磅
    if size_pt is None or abs(size_pt - 24) > 0.01:
        reasons.append(f"字号非 24 磅（实际 {size_pt}）")
    # (6) 加粗
    if not bold:
        reasons.append("未加粗")
    # (4) 深绿色：G 通道占优 且 整体偏暗（亮度 < ~120）
    if not color_hex:
        reasons.append("未识别到字体颜色")
    else:
        r, g, b = hex_to_rgb(color_hex)
        brightness = (r + g + b) / 3
        is_dark_green = (g > r) and (g >= b - 30) and brightness < 130
        if not is_dark_green:
            reasons.append(f"颜色非深绿色（#{color_hex}）")

    if reasons:
        return False, "; ".join(reasons)

    return True, (
        f"第 3 张岗位卡片(右侧)内 '排班岗' 标题：字体={name}, "
        f"字号={size_pt}磅, 加粗={bold}, 颜色=#{color_hex}（深绿）"
    )


def rule_5_rotation(prs):
    """+3  第5页全部卡片和右侧图片：旋转角度调整为 0 度或接近 0 度，
         卡片边线、文字基线和图片上下边缘保持水平。逐点校验：
      (1) 第 5 页"全部卡片"旋转角 ≈ 0°
      (2) 第 5 页"右侧图片"旋转角 ≈ 0°
      (3) 卡片边线水平（即矩形未发生旋转/倾斜）
      (4) 卡片内文字基线水平（每个 run 的字符方向无旋转/无倾斜变形）
      (5) 图片上下边缘水平（picture 未旋转，且未应用任何 a:xfrm rot 或非零 skew）
    """
    s5 = prs.slides[4]
    sw = prs.slide_width

    TOL_DEG = 2.0  # "接近 0 度" 容差

    def _shape_rotation(shp):
        """归一化到 [0, 180]，便于做 ≈0 判断（179° 与 -1° 都视为 1°）。"""
        rot = (getattr(shp, "rotation", 0) or 0) % 360
        if rot > 180:
            rot = 360 - rot
        return rot

    def _shape_skew(shp):
        """读取 spPr/xfrm 上是否存在非零 skew（极少见，保险起见检查）。"""
        x = shp._element.find(".//a:xfrm", NS)
        if x is None:
            return 0
        # rot 已通过 shp.rotation 拿到；这里检查 flipH/flipV 不影响"水平"
        # OOXML 标准里 a:xfrm 没有显式 skew 字段，倾斜通常通过 rot 表达；
        # 若存在 a:custGeom 旋转节点也已经被 rotation 涵盖
        return 0

    # ---- 卡片识别：第 5 页中宽 ≥ 2″、高 ≥ 1″ 的 AUTO_SHAPE 视作卡片容器 ----
    cards = []
    for shp in s5.shapes:
        if shp.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        w = shp.width or 0
        h = shp.height or 0
        if w >= Emu(2 * 914400) and h >= Emu(1 * 914400):
            cards.append(shp)

    if not cards:
        return False, "第 5 页未识别到任何卡片"

    # ---- 右侧图片识别：shape_type == PICTURE 且中心 x ≥ slide_width / 2 ----
    right_pics = [
        shp for shp in s5.shapes
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE
        and ((shp.left or 0) + (shp.width or 0) / 2) >= sw * 0.5
    ]
    if not right_pics:
        return False, "第 5 页未识别到右侧图片"

    bad = []

    # ---- (1)(3) 全部卡片旋转 ≈ 0° → 卡片边线水平 ----
    for c in cards:
        rot = _shape_rotation(c)
        if rot > TOL_DEG:
            bad.append(f"卡片 {c.name} 旋转 {rot:.1f}°（边线不水平）")
        if _shape_skew(c):
            bad.append(f"卡片 {c.name} 存在倾斜变形")

    # ---- (4) 卡片内文字基线水平：文本框自身及其每个段落必须无旋转 ----
    for shp in s5.shapes:
        if not shp.has_text_frame:
            continue
        # 文本框本体旋转
        rot = _shape_rotation(shp)
        if rot > TOL_DEG:
            bad.append(f"文字框 {shp.name} 旋转 {rot:.1f}°（基线不水平）")
            continue
        # 段落内 run 是否带 vert/rot（极少见的字符级旋转/竖排）
        bodyPr = shp._element.find(".//a:bodyPr", NS)
        if bodyPr is not None:
            vert = bodyPr.get("vert", "")
            if vert and vert not in ("horz", ""):
                bad.append(f"文字框 {shp.name} 文本方向 vert={vert}（基线不水平）")
            body_rot = bodyPr.get("rot")
            if body_rot:
                try:
                    deg = int(body_rot) / 60000.0
                    if abs(deg) > TOL_DEG:
                        bad.append(f"文字框 {shp.name} bodyPr rot={deg:.1f}°")
                except ValueError:
                    pass

    # ---- (2)(5) 右侧图片旋转 ≈ 0°：上下边缘水平 ----
    for pic in right_pics:
        rot = _shape_rotation(pic)
        if rot > TOL_DEG:
            bad.append(f"右侧图片 {pic.name} 旋转 {rot:.1f}°（上下边缘不水平）")
        # picture 内部 blipFill / xfrm 的 rot
        pic_xfrm = pic._element.find(".//p:spPr/a:xfrm", NS)
        if pic_xfrm is not None:
            r = pic_xfrm.get("rot")
            if r:
                try:
                    deg = int(r) / 60000.0
                    if abs(deg) > TOL_DEG:
                        bad.append(f"右侧图片 {pic.name} xfrm rot={deg:.1f}°")
                except ValueError:
                    pass

    if bad:
        return False, "; ".join(bad)

    return True, (
        f"第 5 页共 {len(cards)} 张卡片、{len(right_pics)} 张右侧图片：旋转角均 ≤ {TOL_DEG}°，"
        f"卡片边线/文字基线/图片上下边缘均水平"
    )


def rule_5_background(prs):
    """+3  第5页背景更换为与第2、3页一致或高度接近的浅蓝色背景，
         不再使用米黄色背景。逐点校验：
      (1) 能读取到第 5 页的背景颜色
      (2) 第 5 页背景是"浅蓝色"
      (3) 与第 2、3 页背景"一致或高度接近"
      (4) 不再是"米黄色"背景
    任一不满足即不加分。

    严格化（回应反馈）：
      · 背景色不再只读 p:cSld/p:bg 下的 a:srgbClr；改为 resolve_slide_background_hex
        依次解析 slide bg → 全页背景形状 → layout bg → master bg，并支持
        srgbClr/schemeClr/sysClr/prstClr + lumMod/lumOff/shade/tint 主题色链。
    """
    bg5 = resolve_slide_background_hex(prs, prs.slides[4])
    bg2 = resolve_slide_background_hex(prs, prs.slides[1])
    bg3 = resolve_slide_background_hex(prs, prs.slides[2])

    # ---- (1) 必须能读取到第 5 页背景 ----
    if bg5 is None:
        return False, "无法读取第 5 页背景色（已尝试 slide/形状/layout/master）"

    # ---- (4) 不再是米黄色 ----
    if is_beige(bg5):
        return False, f"第 5 页仍为米黄色背景 #{bg5}"

    # ---- (2) 浅蓝色 ----
    if not is_blueish(bg5):
        return False, f"第 5 页背景 #{bg5} 不在浅蓝色范围"

    # ---- (3) 与第 2、3 页一致 或 高度接近 ----
    # "一致"：完全相等；"高度接近"：与第 2 页或第 3 页其中之一颜色欧氏距离 ≤ 40
    if bg2 is None and bg3 is None:
        return False, "无法读取第 2、3 页背景色，无法比较"

    same_as_2 = (bg2 is not None and bg5 == bg2)
    same_as_3 = (bg3 is not None and bg5 == bg3)
    d2 = color_distance(bg5, bg2) if bg2 else 9999
    d3 = color_distance(bg5, bg3) if bg3 else 9999
    closeness = min(d2, d3)

    if not (same_as_2 or same_as_3 or closeness <= 40):
        return False, (
            f"第 5 页背景 #{bg5} 与第 2/3 页 (#{bg2}/#{bg3}) 既不一致也不接近 "
            f"(最小色距 {closeness:.1f} > 40)"
        )

    relation = "一致" if (same_as_2 or same_as_3) else f"高度接近(最小色距 {closeness:.1f})"
    return True, (
        f"第 5 页背景 #{bg5}（浅蓝、非米黄），与第 2/3 页 #{bg2}/#{bg3} {relation}"
    )


def rule_6_optimize(prs):
    """+1  第6页"优化动作"区域："将敏感表达改为场景化、说明式表达" 内容
         出现次数不能超过 1。逐点校验：
      (1) 第 6 页存在 "优化动作" 区域（含标题文本框）
      (2) 在该区域的内容框中，目标语句 "将敏感表达改为场景化、说明式表达"
          的出现次数 ≤ 1（0 次或 1 次都算通过；≥ 2 次则不通过）
    """
    s6 = prs.slides[5]

    # ---- (1) 定位 "优化动作" 标题 ----
    title_shp = None
    for shp in s6.shapes:
        if shp.has_text_frame and shp.text_frame.text.strip() == "优化动作":
            title_shp = shp
            break
    if title_shp is None:
        return False, "第 6 页未找到 '优化动作' 区域标题"

    # 在标题正下方、横向重叠的文本框中找到"优化动作"区域的内容框
    title_cx = (title_shp.left or 0) + (title_shp.width or 0) / 2
    title_w = title_shp.width or 0
    title_cy = (title_shp.top or 0) + (title_shp.height or 0) / 2

    content_shp = None
    for shp in s6.shapes:
        if shp is title_shp or not shp.has_text_frame:
            continue
        if shp.text_frame.text.strip() == "优化动作":
            continue
        if not shp.text_frame.text.strip():
            continue
        cx = (shp.left or 0) + (shp.width or 0) / 2
        cy = (shp.top or 0) + (shp.height or 0) / 2
        # 横向中心在标题宽度范围内 + 纵向中心在标题中心下方
        if abs(cx - title_cx) <= max(title_w, 1) / 2 + 1 and cy > title_cy:
            content_shp = shp
            break

    if content_shp is None:
        return False, "第 6 页未找到 '优化动作' 区域的内容框"

    # ---- (2) 目标语句出现次数 ≤ 1 ----
    target = "将敏感表达改为场景化、说明式表达"
    # 以"段落"为粒度统计：每个段落若包含目标语句即计 1 次
    # （去除行首项目符号 • · - 后再判断，避免误计）
    occurrences = 0
    for p in content_shp.text_frame.paragraphs:
        line = re.sub(r"^[•·\-\s]+", "", p.text.strip())
        if target in line:
            occurrences += 1

    if occurrences > 1:
        return False, (
            f"'优化动作' 区域内 '{target}' 出现 {occurrences} 次（超过 1 次）"
        )

    return True, (
        f"第 6 页 '优化动作' 区域内 '{target}' 出现 {occurrences} 次（≤ 1，通过）"
    )


def rule_9_two_blocks(prs):
    """+5  第9页 PPT 除标题和背景外只出现两个板块："近期" 板块和 "长期" 板块，
         两个板块居中放置，板块下方无其余内容。逐点校验：
      (1) 除标题(本页大标题)和背景外，板块标题集合恰好为 {"近期", "长期"}
          —— 不能出现 "中期" / "短期" 等其他板块
      (2) 除标题和背景外，没有"无关"的其他可见内容形状
      (3) 两个板块居中放置（两板块整体在页面水平方向居中、左右对称）
      (4) 两个板块下方无其余内容
    """
    s9 = prs.slides[8]
    sw = prs.slide_width
    sh = prs.slide_height

    # ---- 识别本页"大标题"（页面顶部、含本页标题文本）----
    # 第 9 页大标题为 "后续规划"。为通用起见：取 top < 页高 15% 且文字非空的最顶端文本框
    title_shapes = []
    for shp in s9.shapes:
        if not shp.has_text_frame:
            continue
        if not shp.text_frame.text.strip():
            continue
        if (shp.top or 0) < sh * 0.15:
            title_shapes.append(shp)
    title_top_y = min(((s.top or 0) + (s.height or 0)) for s in title_shapes) if title_shapes else 0

    # ---- 标识"背景元素"：占满或几乎占满整页的 PICTURE / AUTO_SHAPE 视为背景，跳过 ----
    def _is_background(shp):
        w = shp.width or 0
        h = shp.height or 0
        area_ratio = (w * h) / float(sw * sh)
        return area_ratio >= 0.85

    # ---- 收集"非标题、非背景、非空"的内容 shape ----
    BLOCK_TITLES = {"近期", "中期", "长期", "短期"}
    block_title_shapes = []   # 板块标题文本框
    block_body_shapes = []    # 板块内容文本框（位置紧邻某个板块标题下方）
    other_shapes = []         # 其他与上述均无关的可见内容

    # 先确定板块标题
    for shp in s9.shapes:
        if shp in title_shapes:
            continue
        if _is_background(shp):
            continue
        if shp.has_text_frame and shp.text_frame.text.strip() in BLOCK_TITLES:
            block_title_shapes.append(shp)

    # ---- (1) 板块标题集合必须恰好 == {"近期", "长期"} ----
    names = sorted({s.text_frame.text.strip() for s in block_title_shapes})
    if names != ["近期", "长期"]:
        return False, f"第 9 页板块标题为 {names}，不是恰好 '近期'+'长期'"

    # 找出每个板块标题对应的"内容文本框"（横向中心与标题重叠、纵向在标题之下）
    for shp in s9.shapes:
        if shp in title_shapes or shp in block_title_shapes:
            continue
        if _is_background(shp):
            continue
        if not shp.has_text_frame or not shp.text_frame.text.strip():
            # 装饰性形状（线条/小圆点等），不视为"其余内容"——其没有可见文字
            # 但仍要校验其纵向位置（细则要求板块下方无其余内容）
            other_shapes.append(shp)
            continue
        belongs = False
        for tshp in block_title_shapes:
            tcx = (tshp.left or 0) + (tshp.width or 0) / 2
            tw = tshp.width or 0
            cx = (shp.left or 0) + (shp.width or 0) / 2
            if abs(cx - tcx) <= max(tw, 1) / 2 + 1 and (shp.top or 0) > (tshp.top or 0):
                block_body_shapes.append(shp)
                belongs = True
                break
        if not belongs:
            other_shapes.append(shp)

    # ---- (2) 不能出现与板块/标题/背景无关的额外文字内容 ----
    extra_texts = [
        s.text_frame.text.strip() for s in other_shapes
        if s.has_text_frame and s.text_frame.text.strip()
    ]
    if extra_texts:
        return False, f"第 9 页除标题/背景/板块外仍有其他文字内容：{extra_texts}"

    # ---- (4) 两个板块下方无其余内容 ----
    # 取两个板块（标题 + 其内容框）的最大底边 y，再检查 other_shapes 中是否有位置低于它
    block_bottom = 0
    for shp in (block_title_shapes + block_body_shapes):
        block_bottom = max(block_bottom, (shp.top or 0) + (shp.height or 0))

    below_items = []
    for shp in other_shapes:
        if (shp.top or 0) >= block_bottom - 1:
            # 形状位于板块下方
            below_items.append(f"{shp.name} top={emu_to_cm(shp.top or 0):.1f}cm")

    if below_items:
        return False, f"第 9 页板块下方仍有其余内容：{below_items}"

    # ---- (3) 两个板块居中放置：两个板块的整体中心 x ≈ 页面中点；左右对称 ----
    centers = []
    for tshp in block_title_shapes:
        # 板块以"标题 + 内容框"整体作为一个板块的水平包围盒
        body = next(
            (b for b in block_body_shapes if abs(((b.left or 0) + (b.width or 0) / 2) -
                                                 ((tshp.left or 0) + (tshp.width or 0) / 2))
             <= max(tshp.width or 0, 1) / 2 + 1),
            None,
        )
        l = (tshp.left or 0)
        r = l + (tshp.width or 0)
        if body is not None:
            l = min(l, body.left or 0)
            r = max(r, (body.left or 0) + (body.width or 0))
        centers.append((l + r) / 2)

    centers.sort()
    mid = sw / 2
    overall_center = sum(centers) / len(centers)
    left_off = mid - centers[0]
    right_off = centers[1] - mid
    # 容差：±8% 页面宽度
    tol = sw * 0.08
    if abs(overall_center - mid) > tol:
        return False, f"两个板块整体未居中：整体中心 {emu_to_cm(overall_center):.1f}cm，页面中点 {emu_to_cm(mid):.1f}cm"
    if abs(left_off - right_off) > tol:
        return False, (
            f"两个板块未居中对称：左板块距中点 {emu_to_cm(left_off):.1f}cm，"
            f"右板块距中点 {emu_to_cm(right_off):.1f}cm"
        )

    return True, (
        f"第 9 页除标题/背景外仅 '近期'+'长期' 两板块，两板块整体居中且左右对称，"
        f"板块下方无其余内容"
    )


def rule_9_long_term(prs):
    """+1  第9页右侧 "长期" 板块：三行内容为
         "搭建标准化复盘范式" "强化岗位协同配合效能" "赋能团队稳步长效发展"，
         每句独立一行。逐点校验：
      (1) 位于第 9 页
      (2) "长期" 板块位于页面右侧
      (3) 该板块包含三行内容，依次为指定三句
      (4) 每句独立成一行（一个段落只对应一句）
    """
    s9 = prs.slides[8]
    sw = prs.slide_width

    # ---- (1)(2) 定位"长期"标题，必须位于页面右侧 ----
    long_title = None
    for shp in s9.shapes:
        if shp.has_text_frame and shp.text_frame.text.strip() == "长期":
            long_title = shp
            break
    if long_title is None:
        return False, "第 9 页未找到 '长期' 标题"

    title_cx = (long_title.left or 0) + (long_title.width or 0) / 2
    if title_cx < sw * 0.5:
        return False, f"'长期' 板块位于左侧（中心 x={emu_to_cm(title_cx):.1f}cm），不在右侧"

    # ---- 定位"长期"板块下方的内容框 ----
    content_shp = None
    for shp in s9.shapes:
        if shp is long_title or not shp.has_text_frame:
            continue
        if not shp.text_frame.text.strip():
            continue
        cx = (shp.left or 0) + (shp.width or 0) / 2
        if abs(cx - title_cx) <= max(long_title.width or 0, 1) / 2 + 1 \
                and (shp.top or 0) > (long_title.top or 0):
            content_shp = shp
            break

    if content_shp is None:
        return False, "未在 '长期' 板块下方找到内容框"

    # ---- (3)(4) 三行内容、每句独立一行：以"段落"为粒度比较 ----
    expected = [
        "搭建标准化复盘范式",
        "强化岗位协同配合效能",
        "赋能团队稳步长效发展",
    ]
    paras = [p.text.strip() for p in content_shp.text_frame.paragraphs if p.text.strip()]

    if len(paras) != 3:
        return False, f"'长期' 板块行数为 {len(paras)}（需恰好 3 行）：{paras}"

    if paras != expected:
        return False, f"'长期' 板块内容为 {paras}，与期望 {expected} 不一致"

    return True, "第 9 页右侧 '长期' 板块三行内容完全匹配且每句独立成行"


def rule_10_chart_title(prs):
    """+1  第10页曲线图标题：使用 "月度营收与盈利比值趋势" 或语义一致的标题。
         逐点校验：
      (1) 第 10 页存在曲线图（图表对象）
      (2) 该图表"有标题"
      (3) 标题文本为 "月度营收与盈利比值趋势"，或语义一致
         （语义一致 ≈ 同时包含 "营收"+"盈利" 关键语义，且含 "趋势/月度/月" 时间维度）
    """
    chart = get_chart_on_slide(prs.slides[9])

    # ---- (1) 必须有图表对象 ----
    if chart is None:
        return False, "第 10 页无图表对象"

    # ---- (2) 必须有标题 ----
    if not chart.has_title:
        return False, "图表无标题"

    title_tf = chart.chart_title.text_frame
    title_text = title_tf.text.strip()
    if not title_text:
        return False, "标题为空白"

    # ---- (3) 标题文本：完全匹配 或 语义一致 ----
    expected = "月度营收与盈利比值趋势"
    semantic_ok = (
        ("营收" in title_text)
        and ("盈利" in title_text)
        and (("趋势" in title_text) or ("月度" in title_text) or ("月" in title_text))
    )
    if title_text != expected and not semantic_ok:
        return False, f"标题 '{title_text}' 与 '{expected}' 不一致也不语义相近"

    match_kind = "完全匹配" if title_text == expected else "语义一致"
    return True, f"第 10 页图表标题 '{title_text}'（{match_kind}）"


def rule_10_chart_size(prs):
    """+1  第10页曲线图版式：宽度约13至17厘米、高度约7至10厘米，
         与结束页其他元素之间保留合理间距。逐点校验：
      (1) 第 10 页存在曲线图对象
      (2) 宽度约 13~17 cm
      (3) 高度约 7~10 cm
      (4) 与结束页其他元素之间保留合理间距（不与其他元素重叠/紧贴）
    """
    # ---- (1) 必须存在曲线图对象 ----
    chart_shape = None
    for shp in prs.slides[9].shapes:
        if getattr(shp, "has_chart", False):
            chart_shape = shp
            break
    if chart_shape is None:
        return False, "第 10 页无图表对象"

    # ---- (2) 宽度约 13~17 cm ----
    w_cm = emu_to_cm(chart_shape.width)
    if not (13 <= w_cm <= 17):
        return False, f"图表宽度 {w_cm:.2f}cm 不在 13~17cm 范围"

    # ---- (3) 高度约 7~10 cm ----
    h_cm = emu_to_cm(chart_shape.height)
    if not (7 <= h_cm <= 10):
        return False, f"图表高度 {h_cm:.2f}cm 不在 7~10cm 范围"

    # ---- (4) 与结束页其他元素之间保留合理间距 ----
    # "合理间距" 判定：图表与"真实需避让元素"（标题、姓名、正文、装饰对象等）
    # 不相交，且最小矩形距离 ≥ MIN_GAP_CM。
    #
    # 严格化（回应反馈）：不再遍历所有非零尺寸 shape，会排除：
    #   ① 覆盖近整页的背景形状 / 装饰底图（面积 ≥ 页面 90%）
    #   ② 母版/版式层继承的背景（slide.shapes 遍历的是 slide 级，不含 layout/master；
    #      python-pptx 的 slide.shapes 已经天然排除 layout/master 形状，这里再做一层
    #      面积兜底，避免"作者在 slide 上放了一张全屏背景图"造成误判）
    #   ③ 图表内部/子元素（chart_shape 自身及其 group 子形状）
    #   ④ 完全无内容的空文本框（无文字且无实心填充/边线的占位）
    MIN_GAP_CM = 0.2

    def _rect_of(shp):
        return (
            shp.left or 0,
            shp.top or 0,
            (shp.left or 0) + (shp.width or 0),
            (shp.top or 0) + (shp.height or 0),
        )

    slide10 = prs.slides[9]
    sw = prs.slide_width
    sh = prs.slide_height
    page_area = max(sw * sh, 1)

    # 图表自身（+group 子孙）需被排除
    def _iter_self_and_descendants(root):
        stack = [root]
        while stack:
            cur = stack.pop()
            yield cur
            if getattr(cur, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                try:
                    stack.extend(list(cur.shapes))
                except Exception:
                    pass

    chart_self_ids = {id(s) for s in _iter_self_and_descendants(chart_shape)}

    def _looks_like_background(shp):
        w = shp.width or 0
        h = shp.height or 0
        if w * h >= page_area * 0.9:
            return True
        # 图片且尺寸接近整页短边任一维度、且没有文本 → 装饰底图
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE and (w >= sw * 0.95 or h >= sh * 0.95):
            return True
        return False

    def _is_empty_placeholder(shp):
        # 无文本 + 非图片 + 无实心填充 + 无可见线条 → 视为空占位，不作为需避让元素
        try:
            if shp.has_text_frame and shp.text_frame.text.strip():
                return False
        except Exception:
            return False
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return False
        try:
            sp_pr = shp._element.find(".//p:spPr", NS)
        except Exception:
            sp_pr = None
        if sp_pr is not None:
            if sp_pr.find("a:solidFill", NS) is not None:
                return False
            if sp_pr.find("a:ln//a:solidFill", NS) is not None:
                return False
        return True

    cr = _rect_of(chart_shape)
    min_gap_cm = None
    overlapping = []
    scanned = 0
    for shp in slide10.shapes:
        if id(shp) in chart_self_ids:
            continue
        if (shp.width or 0) == 0 or (shp.height or 0) == 0:
            continue
        if _looks_like_background(shp):
            continue
        if _is_empty_placeholder(shp):
            continue

        scanned += 1
        r = _rect_of(shp)
        dx = max(0, max(cr[0] - r[2], r[0] - cr[2]))
        dy = max(0, max(cr[1] - r[3], r[1] - cr[3]))
        if dx == 0 and dy == 0:
            overlapping.append(shp.name)
            continue
        gap_cm = emu_to_cm((dx ** 2 + dy ** 2) ** 0.5)
        if min_gap_cm is None or gap_cm < min_gap_cm:
            min_gap_cm = gap_cm

    if overlapping:
        return False, f"图表与真实需避让元素重叠：{overlapping}"
    if min_gap_cm is not None and min_gap_cm < MIN_GAP_CM:
        return False, f"图表与最近需避让元素间距仅 {min_gap_cm:.2f}cm（需 ≥ {MIN_GAP_CM}cm）"

    gap_desc = f"{min_gap_cm:.2f}cm" if min_gap_cm is not None else f"页内无真实需避让元素(扫描 {scanned})"
    return True, (
        f"第 10 页图表 {w_cm:.1f}×{h_cm:.1f}cm（宽 13~17、高 7~10），"
        f"与真实需避让元素最小间距 {gap_desc}（合理，已排除背景/图表自身/空占位）"
    )


# --------------------------- 维度2 扣分项 ---------------------------


# --------------------------- 主程序 ---------------------------


SCRIPT_ID = "063"
MAX_SCORE = 5 + 1 + 3 + 3 + 1 + 5 + 1 + 1 + 1  # 维度二所有加分项满分之和 = 21


def _find_pptx_in_dir(dir_path):
    """在指定目录内定位待评估的 .pptx 文件，返回绝对路径。"""
    if not os.path.isdir(dir_path):
        raise FileNotFoundError(f"目录不存在或不是目录：{dir_path}")
    candidates = [
        name for name in os.listdir(dir_path)
        if name.lower().endswith(".pptx")
        and not name.startswith("~$")  # 排除 Office 临时锁文件
    ]
    if not candidates:
        raise FileNotFoundError(f"目录 {dir_path} 下未找到 .pptx 文件")
    # 若有多份，优先选择非临时文件中最新修改的一份，避免误挑
    candidates.sort(
        key=lambda n: os.path.getmtime(os.path.join(dir_path, n)),
        reverse=True,
    )
    return os.path.join(dir_path, candidates[0])


def evaluate(dir_path: str) -> dict:
    """脚本对外唯一入口。

    参数：
        dir_path: 脚本所在目录路径；脚本自行在该目录内定位 .pptx 文档。

    返回：
        结构化 dict，字段含义详见 §2.2 统一约定。
    """
    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": MAX_SCORE,
    }

    try:
        pptx_path = _find_pptx_in_dir(dir_path)
        result["file_name"] = os.path.basename(pptx_path)

        ok, fails = check_dim1(pptx_path)
        if not ok:
            result["dim1_pass"] = False
            result["dim1_reason"] = "；".join(fails)
            result["total_score"] = 0
            return result

        result["dim1_pass"] = True

        prs = Presentation(pptx_path)

        # 评分细则原文（来自打分细则）
        rules = [
            ("第10页曲线图对象：在结束页右侧或中部空白区域，插入可编辑的折线图或平滑曲线图，"
             "包含\"营收数值\"和\"盈利比值\"两个纵轴数据系列，两条曲线颜色不同。"
             "横轴从左至右依次显示\"1月、2月、3月、4月、5月\"，月份顺序正确。"
             "营收数值系列：五个月数据依次为68000、352000、486000、423000、579000。"
             "盈利比值系列：五个月数据依次为1.42、1.56、1.71、1.63、1.85。"
             "准确显示\"营收数值\"\"盈利比值\"，图例颜色与两条曲线一一对应。", 5, rule_10_chart),
            ("第2页右侧第3个岗位卡片：空白处有标题\"排班岗\"字体为深绿色微软雅黑24磅加粗",
             1, rule_2_paiban),
            ("第5页全部卡片和右侧图片：旋转角度调整为0度或接近0度，"
             "卡片边线、文字基线和图片上下边缘保持水平。", 3, rule_5_rotation),
            ("第5页背景更换为与第2、3页一致或高度接近的浅蓝色背景，不再使用米黄色背景。",
             3, rule_5_background),
            ("第6页\"优化动作\"区域：\"将敏感表达改为场景化、说明式表达\"内容出现次数不能超过1",
             1, rule_6_optimize),
            ("第9页PPT除标题和背景外只出现两个板块：\"近期\"板块和\"长期\"板块，"
             "两个板块居中放置，板块下方无其余内容", 5, rule_9_two_blocks),
            ("第9页右侧\"长期\"板块：三行内容为\"搭建标准化复盘范式\"\"强化岗位协同配合效能\""
             "\"赋能团队稳步长效发展\"，每句独立一行。", 1, rule_9_long_term),
            ("第10页曲线图标题：使用\"月度营收与盈利比值趋势\"或语义一致的标题。",
             1, rule_10_chart_title),
            ("第10页曲线图版式：宽度约13至17厘米、高度约7至10厘米，"
             "与结束页其他元素之间保留合理间距。", 1, rule_10_chart_size),
        ]

        deductions = []

        total = 0
        items = []

        for content, score, fn in rules:
            try:
                ok_r, _ = fn(prs)
            except Exception:
                ok_r = False
            delta = score if ok_r else 0
            total += delta
            items.append({
                "rule": content,
                "max_delta": score,
                "delta": delta,
                "hit": bool(ok_r),
                "detail": "",
            })

        for content, score, fn in deductions:
            try:
                ok_r, _ = fn(prs)
            except Exception:
                ok_r = False
            delta = score if ok_r else 0
            total += delta
            items.append({
                "rule": content,
                "max_delta": score,
                "delta": delta,
                "hit": bool(ok_r),
                "detail": "",
            })

        result["dim2_items"] = items
        result["total_score"] = total
        return result

    except Exception as e:
        result["status"] = "error"
        result["error"] = f"{type(e).__name__}: {e}"
        result["total_score"] = 0
        return result


if __name__ == "__main__":
    # 本地调试入口：默认使用脚本所在目录，也可通过命令行覆盖
    default_dir = os.path.dirname(os.path.abspath(__file__))
    target_dir = sys.argv[1] if len(sys.argv) > 1 else default_dir
    print(json.dumps(evaluate(target_dir), ensure_ascii=True, indent=2))
