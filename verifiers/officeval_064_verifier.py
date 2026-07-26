# -*- coding: utf-8 -*-
"""
自动评估 “核心技术_照片排版版.pptx”
评分逻辑：
  - 维度1（可用与可修改性）不通过 → 直接 0 分；
  - 维度1 通过 → 累计维度2 的得分点 / 扣分点。
"""
import os
import sys
import json
import zipfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

SCRIPT_ID = "064"
EMU_PER_CM = 360000.0


def cm(emu):
    return (emu or 0) / EMU_PER_CM


def in_range(v, lo, hi, tol=0.05):
    """带一点点容差地判断 v∈[lo,hi]。"""
    return (lo - tol) <= v <= (hi + tol)


def in_range_strict(v, lo, hi):
    """严格按细则区间判断，不带容差。"""
    return lo <= v <= hi


def rect(shape):
    return (cm(shape.left), cm(shape.top), cm(shape.width), cm(shape.height))


def _center_inside(shape, box_rect) -> bool:
    """判断 shape 的中心点是否落在矩形 box_rect=(left, top, width, height) 内。
    坐标单位均为 cm。允许 0.2cm 的容差处理精度问题。
    """
    cx = cm(shape.left) + cm(shape.width) / 2
    cy = cm(shape.top) + cm(shape.height) / 2
    bl, bt, bw, bh = box_rect
    tol = 0.2
    return (bl - tol) <= cx <= (bl + bw + tol) and (bt - tol) <= cy <= (bt + bh + tol)


def _locate_module_box(slide, keyword: str, fallback_rect):
    """按模块标题文本定位模块内容框（背景矩形）。

    策略：
      1. 在页面上找到含 `keyword` 的文本框 title_shape；
      2. 在同页所有非图片/非文本的大矩形形状里，挑一个"包含 title_shape 中心点、
         宽度 ≥ 10cm、高度 ≥ 5cm" 的作为背景框；
      3. 兜底：若上面两步任何一步失败，返回 `fallback_rect`。

    返回 (left, top, width, height)，单位 cm。
    """
    title_shape = None
    for sh in slide.shapes:
        if not getattr(sh, "has_text_frame", False):
            continue
        try:
            txt = sh.text_frame.text or ""
        except Exception:
            txt = ""
        if keyword and keyword in txt:
            title_shape = sh
            break
    if title_shape is None:
        return fallback_rect

    tx = cm(title_shape.left) + cm(title_shape.width) / 2
    ty = cm(title_shape.top) + cm(title_shape.height) / 2

    best = None
    best_area = 0.0
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        # 排除承载文字的形状本身
        try:
            if getattr(sh, "has_text_frame", False) and (sh.text_frame.text or "").strip():
                continue
        except Exception:
            pass
        w, h = cm(sh.width), cm(sh.height)
        if w < 10 or h < 5:
            continue
        left, top = cm(sh.left), cm(sh.top)
        if not (left <= tx <= left + w and top <= ty <= top + h):
            continue
        area = w * h
        # 选包含标题中心且面积最小的那个（更贴合模块框而不是整页背景）
        if best is None or area < best_area:
            best = (left, top, w, h)
            best_area = area
    return best if best is not None else fallback_rect


def overlap(a, b, min_overlap_cm=0.3):
    """两个矩形是否“显著”重叠（>= min_overlap_cm 厘米的交集边长）。"""
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0.0, min(ay + ah, by + by + bh) - max(ay, by))
    # 修正 iy
    iy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
    return ix > min_overlap_cm and iy > min_overlap_cm


# ---------- 维度1 ----------
def check_dimension1(path: str):
    """交付文件为 .pptx 格式，能够正常打开。

    只做两件事：
      1) 后缀是 .pptx；
      2) 能够被 zip / python-pptx 正常打开。
    其它（页数、图片数、遮挡、可编辑、放映等）不再作为维度一的判据。
    """
    reasons = []
    if not path.lower().endswith(".pptx"):
        return False, ["文件不是 .pptx 格式"]
    if not os.path.exists(path):
        return False, ["文件不存在"]
    try:
        with zipfile.ZipFile(path) as zf:
            _ = zf.namelist()
    except Exception as e:
        return False, [f"文件无法打开（zip 损坏）：{e}"]
    try:
        _ = Presentation(path)
    except Exception as e:
        return False, [f"python-pptx 解析失败：{e}"]

    return True, reasons


# ---------- 维度2 工具：定位 4 个内容框中的图片 ----------
def collect_box_pictures(prs, layout_specs):
    """
    返回 {(slide_idx, side): [pictures按规则期望顺序排列]}
      slide_idx ∈ {0,1}; side ∈ {'L','R'} 依据图片左边距是否 < 页宽一半。
    排序：第1页是“顶部大图 → 左下小图 → 右下小图”，按 (top, left) 升序即可；
          第2页是“左上小图 → 左下小图 → 右侧大图”，按 (left, top) 升序更贴近规则。
    """
    boxes = {}
    slide_w = cm(prs.slide_width)
    for si, slide in enumerate(prs.slides):
        pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        local = {}
        for p in pics:
            side = "L" if cm(p.left) + cm(p.width) / 2 < slide_w / 2 else "R"
            local.setdefault((si, side), []).append(p)
        # 第 1 页(si=0)：顶部大图独占一行→按 top 主、left 副；
        # 第 2 页(si=1)：左侧两小图在左、大图在右→按 left 主、top 副。
        for k, v in local.items():
            if si == 0:
                v.sort(key=lambda s: (round(cm(s.top), 1), round(cm(s.left), 1)))
            else:
                v.sort(key=lambda s: (round(cm(s.left), 1), round(cm(s.top), 1)))
            boxes[k] = v
    return boxes


# ---------- 维度2 ----------
def check_dimension2(path):
    prs = Presentation(path)
    points = []  # (delta, label, hit)

    boxes = collect_box_pictures(prs, None)

    # 工具：按尺寸规格校验一个 box 的 3 张图
    def check_box(slide_idx, side, specs, label):
        """
        specs: [(w_lo,w_hi,h_lo,h_hi), ...]  按阅读顺序对应 box 中的图
        """
        pics = boxes.get((slide_idx, side), [])
        if len(pics) != len(specs):
            return False, f"{label}：图片数={len(pics)} 期望 {len(specs)}"
        detail = []
        ok_all = True
        for i, (p, spec) in enumerate(zip(pics, specs)):
            pw, ph = cm(p.width), cm(p.height)
            w_lo, w_hi, h_lo, h_hi = spec
            ok = in_range(pw, w_lo, w_hi) and in_range(ph, h_lo, h_hi)
            ok_all = ok_all and ok
            detail.append(f"图{i+1}({pw:.2f}×{ph:.2f}cm,期望{w_lo}-{w_hi}×{h_lo}-{h_hi}){'✓' if ok else '✗'}")
        return ok_all, label + "：" + " | ".join(detail)

    # +5: 第1页ppt左侧“智能变速总成”框
    # 细则三条子项：
    #   ① 框顶部放置 宽 13.3-13.5cm、高 2.9-3.1cm 的图片
    #   ② 大图下方左侧小图 宽 6.5-6.7cm、高 2.7-2.9cm
    #   ③ 大图下方右侧小图 宽 6.5-6.6cm、高 2.8-2.9cm
    #
    # 严格化（回应反馈）：
    #   · 不能只按"左半页"筛图片；必须先定位"智能变速总成"内容框：
    #       (a) 优先按模块标题文本"智能变速总成"所在文本框位置回溯到其父背景形状；
    #       (b) 兜底：使用 baseline "左侧模块背景@slide1" 的矩形 (1.78, 4.95, 14.78, 11.81)。
    #   · 三张图片必须其中心点落在该内容框内；框内数量 ≠ 3 视为不合格。
    #   · 顶部大图 top 必须显著在小图之上；两小图分别在左下/右下位置。
    def check_slide1_left():
        details = []
        slide = prs.slides[0]

        module_rect = _locate_module_box(
            slide,
            keyword="智能变速总成",
            fallback_rect=(1.78, 4.95, 14.78, 11.81),
        )
        details.append(
            f"模块框=(L={module_rect[0]:.2f}, T={module_rect[1]:.2f}, "
            + f"W={module_rect[2]:.2f}, H={module_rect[3]:.2f})cm"
        )

        pics_in_box = [
            p for p in slide.shapes
            if p.shape_type == MSO_SHAPE_TYPE.PICTURE
            and _center_inside(p, module_rect)
        ]
        if len(pics_in_box) != 3:
            details.append(f"框内图片数={len(pics_in_box)}（期望 3）")
            return False, "第1页左·智能变速总成：" + " | ".join(details)

        # 顶部大图 = top 最小；两小图 = 其余两张按 left 排序
        top_pic = min(pics_in_box, key=lambda p: cm(p.top))
        below = [p for p in pics_in_box if p is not top_pic]
        below.sort(key=lambda p: cm(p.left))
        left_small, right_small = below[0], below[1]

        # 位置语义：顶部大图应显著在两小图上方；两小图应位于下半区并左右分布。
        box_left, box_top, box_w, box_h = module_rect
        vert_split = box_top + box_h * 0.45   # 大约划分上/下两区
        horiz_mid = box_left + box_w / 2
        pos_ok = (
            cm(top_pic.top) < vert_split
            and cm(left_small.top) >= vert_split - 0.3
            and cm(right_small.top) >= vert_split - 0.3
            and cm(left_small.left) + cm(left_small.width) / 2 < horiz_mid
            and cm(right_small.left) + cm(right_small.width) / 2 >= horiz_mid
        )
        if not pos_ok:
            details.append("位置语义✗（顶部大图/左下小图/右下小图 位置不匹配）")

        def chk(pic, w_lo, w_hi, h_lo, h_hi, name):
            pw, ph = cm(pic.width), cm(pic.height)
            ok = in_range_strict(pw, w_lo, w_hi) and in_range_strict(ph, h_lo, h_hi)
            details.append(
                f"{name}({pw:.2f}×{ph:.2f}cm,期望{w_lo}-{w_hi}×{h_lo}-{h_hi})"
                + ('✓' if ok else '✗')
            )
            return ok

        ok1 = chk(top_pic,     13.3, 13.5, 2.9, 3.1, "框顶部图片")
        ok2 = chk(left_small,   6.5,  6.7, 2.7, 2.9, "大图下方左侧小图")
        ok3 = chk(right_small,  6.5,  6.6, 2.8, 2.9, "大图下方右侧小图")
        return (ok1 and ok2 and ok3 and pos_ok), "第1页左·智能变速总成：" + " | ".join(details)

    ok, msg = check_slide1_left()
    points.append((+5,
                   "第1页ppt左侧“智能变速总成”框：框顶部放置对应的宽13.3-13.5cm高2.9-3.1cm的图片；大图下方左侧小图放置对应的宽6.5-6.7cm高2.7-2.9cm的照片；大图下方右侧小图放置对应的宽6.5-6.6cm高2.8-2.9cm的照片",
                   msg, ok))

    # +5: 第1页ppt右侧“高强复合半轴”框
    # 细则三条子项（明确指定为整体第 4、5、6 张图片）：
    #   ① 框顶部 第4张 宽 13.3-13.5cm × 高 2.9-3.1cm
    #   ② 大图下方左侧小图 第5张 宽 6.5-6.7cm × 高 2.7-2.9cm
    #   ③ 大图下方右侧小图 第6张 宽 6.5-6.6cm × 高 2.8-2.9cm
    def check_slide1_right():
        # 严格化（回应反馈）：不再使用 all_pics[3:6] 这种 shape tree 顺序定位。
        # 改为：按"高强复合半轴"标题回溯到内容框（兜底 baseline 右侧模块背景@slide1），
        # 从框内 (center 命中) 的图片按视觉顺序识别：顶部大图 = top 最小；
        # 剩下两张按 left 升序 = 左下 / 右下。
        details = []
        slide = prs.slides[0]

        module_rect = _locate_module_box(
            slide,
            keyword="高强复合半轴",
            fallback_rect=(17.32, 4.95, 14.78, 11.89),
        )
        details.append(
            f"模块框=(L={module_rect[0]:.2f}, T={module_rect[1]:.2f}, "
            + f"W={module_rect[2]:.2f}, H={module_rect[3]:.2f})cm"
        )

        pics_in_box = [
            p for p in slide.shapes
            if p.shape_type == MSO_SHAPE_TYPE.PICTURE
            and _center_inside(p, module_rect)
        ]
        if len(pics_in_box) != 3:
            details.append(f"框内图片数={len(pics_in_box)}（期望 3）")
            return False, "第1页右·高强复合半轴：" + " | ".join(details)

        top_pic = min(pics_in_box, key=lambda p: cm(p.top))
        below = [p for p in pics_in_box if p is not top_pic]
        below.sort(key=lambda p: cm(p.left))
        left_small, right_small = below[0], below[1]

        box_left, box_top, box_w, box_h = module_rect
        vert_split = box_top + box_h * 0.45
        horiz_mid = box_left + box_w / 2
        pos_ok = (
            cm(top_pic.top) < vert_split
            and cm(left_small.top) >= vert_split - 0.3
            and cm(right_small.top) >= vert_split - 0.3
            and cm(left_small.left) + cm(left_small.width) / 2 < horiz_mid
            and cm(right_small.left) + cm(right_small.width) / 2 >= horiz_mid
        )
        if not pos_ok:
            details.append("位置语义✗（顶部/左下/右下 位置不匹配）")

        def chk(pic, w_lo, w_hi, h_lo, h_hi, name):
            pw, ph = cm(pic.width), cm(pic.height)
            ok = in_range_strict(pw, w_lo, w_hi) and in_range_strict(ph, h_lo, h_hi)
            details.append(
                f"{name}({pw:.2f}×{ph:.2f}cm,期望{w_lo}-{w_hi}×{h_lo}-{h_hi})"
                + ('✓' if ok else '✗')
            )
            return ok

        ok1 = chk(top_pic,     13.3, 13.5, 2.9, 3.1, "框顶部图片(视觉第4张)")
        ok2 = chk(left_small,   6.5,  6.7, 2.7, 2.9, "大图下方左侧小图(视觉第5张)")
        ok3 = chk(right_small,  6.5,  6.6, 2.8, 2.9, "大图下方右侧小图(视觉第6张)")
        return (ok1 and ok2 and ok3 and pos_ok), "第1页右·高强复合半轴：" + " | ".join(details)

    ok, msg = check_slide1_right()
    points.append((+5,
                   "第1页ppt右侧“高强复合半轴”框：框顶部放置对应的宽13.3-13.5cm高2.9-3.1cm的第4张图片；大图下方左侧小图放置对应的宽6.5-6.7cm高2.7-2.9cm的第5张照片；大图下方右侧小图放置对应的宽6.5-6.6cm高2.8-2.9cm的第6张照片",
                   msg, ok))

    # +5: 第2页ppt左侧“尾部导流组件”框
    # 细则三条子项：
    #   ① 左上小图 宽 5.1-5.3cm × 高 2.9-3.1cm
    #   ② 左下小图 宽 5.1-5.3cm × 高 2.9-3.1cm
    #   ③ 右侧大图 宽 7.9-8.1cm × 高 6.2-6.4cm
    #
    # 严格化（回应反馈）：
    #   · 不再按"第2页左半页"分组来选图；先按标题"尾部导流组件"回溯到模块背景框
    #     （兜底 baseline "左侧模块背景@slide2" 矩形 (1.78, 5.13, 14.78, 12.01)）；
    #   · 只收集中心点落在模块框内的图片；数量 ≠ 3 直接失败；
    #   · 位置语义按框内坐标区间判定：
    #       左上/左下：中心 x 位于框左半区；右侧大图：中心 x 位于框右半区；
    #       左上 vs 左下：按 top 排序（左上 top 较小）；
    #       右侧大图 top 处于框中部（覆盖左侧两小图高度范围）；
    #   · 通过位置语义后再逐张校验尺寸区间。
    def check_slide2_left():
        details = []
        slide = prs.slides[1]

        module_rect = _locate_module_box(
            slide,
            keyword="尾部导流组件",
            fallback_rect=(1.78, 5.13, 14.78, 12.01),
        )
        details.append(
            f"模块框=(L={module_rect[0]:.2f}, T={module_rect[1]:.2f}, "
            + f"W={module_rect[2]:.2f}, H={module_rect[3]:.2f})cm"
        )

        pics_in_box = [
            p for p in slide.shapes
            if p.shape_type == MSO_SHAPE_TYPE.PICTURE
            and _center_inside(p, module_rect)
        ]
        if len(pics_in_box) != 3:
            details.append(f"框内图片数={len(pics_in_box)}（期望 3）")
            return False, "第2页左·尾部导流组件：" + " | ".join(details)

        box_left, box_top, box_w, box_h = module_rect
        horiz_mid = box_left + box_w / 2

        # 按框内坐标区间划分：左半区两张（左上/左下），右半区一张（右侧大图）
        left_col = [p for p in pics_in_box
                    if cm(p.left) + cm(p.width) / 2 < horiz_mid]
        right_col = [p for p in pics_in_box
                     if cm(p.left) + cm(p.width) / 2 >= horiz_mid]

        if len(left_col) != 2 or len(right_col) != 1:
            details.append(
                f"分列失败：框内左列={len(left_col)}(期望 2)，右列={len(right_col)}(期望 1)"
            )
            return False, "第2页左·尾部导流组件：" + " | ".join(details)

        left_col.sort(key=lambda p: cm(p.top))
        left_top, left_bottom = left_col[0], left_col[1]
        right_big = right_col[0]

        # 位置语义：左上 top 应显著在左下上方；右侧大图 top 应位于框中部
        pos_ok = (
            cm(left_top.top) < cm(left_bottom.top)
            and cm(right_big.top) >= box_top
            and cm(right_big.top) + cm(right_big.height) <= box_top + box_h + 0.5
        )
        if not pos_ok:
            details.append("位置语义✗（左上/左下/右侧大图 位置不匹配）")

        def chk(pic, w_lo, w_hi, h_lo, h_hi, name):
            pw, ph = cm(pic.width), cm(pic.height)
            ok = in_range_strict(pw, w_lo, w_hi) and in_range_strict(ph, h_lo, h_hi)
            details.append(
                f"{name}({pw:.2f}×{ph:.2f}cm,期望{w_lo}-{w_hi}×{h_lo}-{h_hi})"
                + ('✓' if ok else '✗')
            )
            return ok

        ok1 = chk(left_top,    5.1, 5.3, 2.9, 3.1, "左上小图")
        ok2 = chk(left_bottom, 5.1, 5.3, 2.9, 3.1, "左下小图")
        ok3 = chk(right_big,   7.9, 8.1, 6.2, 6.4, "右侧大图")
        return (ok1 and ok2 and ok3 and pos_ok), "第2页左·尾部导流组件：" + " | ".join(details)

    ok, msg = check_slide2_left()
    points.append((+5,
                   "第2页ppt左侧“尾部导流组件”框：左上小图放置对应的宽5.1-5.3cm高2.9-3.1cm的照片；左下小图放置对应的宽5.1-5.3cm高2.9-3.1cm的照片；右侧大图放置对应的宽7.9-8.1cm高6.2-6.4cm的照片",
                   msg, ok))

    # +5: 第2页ppt右侧“座舱轻量布局”框
    # 细则三条子项：
    #   ① 左上小图 宽 5.1-5.3cm × 高 2.9-3.1cm
    #   ② 左下小图 宽 5.1-5.3cm × 高 2.9-3.1cm
    #   ③ 右侧大图 宽 7.9-8.1cm × 高 6.2-6.4cm
    #
    # 严格化（回应反馈）：
    #   · 不再按"第2页右半页"分组来选图；先按标题"座舱轻量布局"回溯到模块背景框
    #     （兜底 baseline "右侧模块背景@slide2" 矩形 (17.32, 5.13, 14.78, 12.03)）；
    #   · 只收集中心点落在模块框内的图片；数量 ≠ 3 直接失败；
    #   · 按框内水平中线分列：左半区两张（左上/左下），右半区一张（右侧大图）；
    #   · 左半区按 top 排序 → 左上 / 左下；右侧大图应位于框中部；
    #   · 通过位置语义后再逐张校验尺寸区间。
    def check_slide2_right():
        details = []
        slide = prs.slides[1]

        module_rect = _locate_module_box(
            slide,
            keyword="座舱轻量布局",
            fallback_rect=(17.32, 5.13, 14.78, 12.03),
        )
        details.append(
            f"模块框=(L={module_rect[0]:.2f}, T={module_rect[1]:.2f}, "
            + f"W={module_rect[2]:.2f}, H={module_rect[3]:.2f})cm"
        )

        pics_in_box = [
            p for p in slide.shapes
            if p.shape_type == MSO_SHAPE_TYPE.PICTURE
            and _center_inside(p, module_rect)
        ]
        if len(pics_in_box) != 3:
            details.append(f"框内图片数={len(pics_in_box)}（期望 3）")
            return False, "第2页右·座舱轻量布局：" + " | ".join(details)

        box_left, box_top, box_w, box_h = module_rect
        horiz_mid = box_left + box_w / 2

        left_col = [p for p in pics_in_box
                    if cm(p.left) + cm(p.width) / 2 < horiz_mid]
        right_col = [p for p in pics_in_box
                     if cm(p.left) + cm(p.width) / 2 >= horiz_mid]

        if len(left_col) != 2 or len(right_col) != 1:
            details.append(
                f"分列失败：框内左列={len(left_col)}(期望 2)，右列={len(right_col)}(期望 1)"
            )
            return False, "第2页右·座舱轻量布局：" + " | ".join(details)

        left_col.sort(key=lambda p: cm(p.top))
        left_top, left_bottom = left_col[0], left_col[1]
        right_big = right_col[0]

        pos_ok = (
            cm(left_top.top) < cm(left_bottom.top)
            and cm(right_big.top) >= box_top
            and cm(right_big.top) + cm(right_big.height) <= box_top + box_h + 0.5
        )
        if not pos_ok:
            details.append("位置语义✗（左上/左下/右侧大图 位置不匹配）")

        def chk(pic, w_lo, w_hi, h_lo, h_hi, name):
            pw, ph = cm(pic.width), cm(pic.height)
            ok = in_range_strict(pw, w_lo, w_hi) and in_range_strict(ph, h_lo, h_hi)
            details.append(
                f"{name}({pw:.2f}×{ph:.2f}cm,期望{w_lo}-{w_hi}×{h_lo}-{h_hi})"
                + ('✓' if ok else '✗')
            )
            return ok

        ok1 = chk(left_top,    5.1, 5.3, 2.9, 3.1, "左上小图")
        ok2 = chk(left_bottom, 5.1, 5.3, 2.9, 3.1, "左下小图")
        ok3 = chk(right_big,   7.9, 8.1, 6.2, 6.4, "右侧大图")
        return (ok1 and ok2 and ok3 and pos_ok), "第2页右·座舱轻量布局：" + " | ".join(details)

    ok, msg = check_slide2_right()
    points.append((+5,
                   "第2页ppt右侧“座舱轻量布局”框：左上小图放置对应的宽5.1-5.3cm高2.9-3.1cm的照片；左下小图放置对应的宽5.1-5.3cm高2.9-3.1cm的照片；右侧大图放置对应的宽7.9-8.1cm高6.2-6.4cm的照片",
                   msg, ok))

    # +3: 第1、2页顶部主标题：外围添加浅蓝色组合边框，边框为可编辑形状。
    #
    # 严格化（回应反馈）：
    #   ① 顶部主标题存在且文本非空；
    #   ② 边框必须是"组合边框"：
    #        · GroupShape 且其内部至少包含 2 个 AUTO_SHAPE/FREEFORM 子对象，或
    #        · 由多个独立 AUTO_SHAPE/FREEFORM 组合成的外围矩形四边（≥2 段）；
    #      且组合外接矩形包围标题矩形。
    #   ③ 边框是"外围线框"而非填充覆盖：
    #        · 每个组成形状必须"有线条颜色 / 线宽 > 0"；
    #        · 若形状为矩形并包围标题，其填充必须是无填充 / 透明（避免覆盖标题）。
    #   ④ 边框颜色为浅蓝色，需解析：直接 srgbClr / schemeClr（读 theme1.xml 的
    #      clrScheme） / prstClr / theme 继承色；带 lumMod/lumOff 修正。
    THEME_SCHEME_CACHE: "dict[str, dict[str, str]]" = {}

    def _load_theme_scheme_from_zip(path: str) -> "dict[str, str]":
        """从 pptx 里读取 ppt/theme/theme1.xml 的 clrScheme。"""
        if path in THEME_SCHEME_CACHE:
            return THEME_SCHEME_CACHE[path]
        out = {}
        try:
            import xml.etree.ElementTree as ET
            with zipfile.ZipFile(path) as zf:
                theme_name = None
                for n in zf.namelist():
                    if n.startswith("ppt/theme/") and n.endswith(".xml"):
                        theme_name = n
                        break
                if theme_name:
                    root = ET.fromstring(zf.read(theme_name))
                    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                    scheme = root.find(f".//{ns_a}themeElements/{ns_a}clrScheme")
                    if scheme is not None:
                        for child in scheme:
                            tag = child.tag.split("}")[-1]
                            srgb = child.find(f"{ns_a}srgbClr")
                            if srgb is not None:
                                v = (srgb.get("val") or "").upper()
                                if len(v) == 6:
                                    out[tag] = v
                                    continue
                            sysc = child.find(f"{ns_a}sysClr")
                            if sysc is not None:
                                v = (sysc.get("lastClr") or "").upper()
                                if len(v) == 6:
                                    out[tag] = v
        except Exception:
            pass
        THEME_SCHEME_CACHE[path] = out
        return out

    theme_scheme = _load_theme_scheme_from_zip(path)

    def _hex_to_rgb(hexv: str):
        try:
            return (int(hexv[0:2], 16), int(hexv[2:4], 16), int(hexv[4:6], 16))
        except Exception:
            return None

    def is_light_blue(rgb):
        if rgb is None:
            return False
        try:
            r, g, b = rgb[0], rgb[1], rgb[2]
        except Exception:
            return False
        return (b > r) and (b > g) and (min(r, g, b) >= 120) and (b >= 180)

    def _apply_lum(hexv: str, lum_mod, lum_off) -> str:
        """近似套用 lumMod/lumOff (val 单位千分数)。"""
        rgb = _hex_to_rgb(hexv)
        if rgb is None:
            return hexv
        r, g, bl = [v / 255.0 for v in rgb]
        mx, mn = max(r, g, bl), min(r, g, bl)
        L = (mx + mn) / 2.0
        new_L = L
        if lum_mod is not None:
            try:
                new_L *= int(lum_mod) / 100000.0
            except Exception:
                pass
        if lum_off is not None:
            try:
                new_L += int(lum_off) / 100000.0
            except Exception:
                pass
        new_L = max(0.0, min(1.0, new_L))
        scale = new_L / L if L > 0 else 1.0
        return "{:02X}{:02X}{:02X}".format(
            min(255, int(r * scale * 255)),
            min(255, int(g * scale * 255)),
            min(255, int(bl * scale * 255)),
        )

    def _resolve_fill_element_color(fill_el) -> "tuple | None":
        """从 <a:solidFill> 或直接颜色元素解析 rgb 三元组，支持 srgb/scheme/prst 与 lumMod/lumOff。"""
        if fill_el is None:
            return None
        import xml.etree.ElementTree as ET  # noqa: F401
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        for tag, resolver in (
            (f"{ns_a}srgbClr", lambda el: el.get("val")),
            (f"{ns_a}schemeClr", lambda el: theme_scheme.get(
                {"tx1": "dk1", "bg1": "lt1", "tx2": "dk2", "bg2": "lt2"}
                .get(el.get("val") or "", el.get("val") or ""))),
            (f"{ns_a}prstClr", lambda el: {
                "lightBlue": "ADD8E6", "skyBlue": "87CEEB",
                "lightSkyBlue": "87CEFA", "powderBlue": "B0E0E6",
                "aliceBlue": "F0F8FF", "cornflowerBlue": "6495ED",
            }.get(el.get("val") or "")),
        ):
            found = fill_el.find(tag) if hasattr(fill_el, "find") else None
            if found is not None:
                hexv = (resolver(found) or "").upper()
                if len(hexv) == 6:
                    lum_mod = None
                    lum_off = None
                    lm = found.find(f"{ns_a}lumMod")
                    lo = found.find(f"{ns_a}lumOff")
                    if lm is not None:
                        lum_mod = lm.get("val")
                    if lo is not None:
                        lum_off = lo.get("val")
                    if lum_mod or lum_off:
                        hexv = _apply_lum(hexv, lum_mod, lum_off)
                    return _hex_to_rgb(hexv)
        return None

    def _shape_line_color_rgb(sh):
        """解析 shape 的 <a:ln>/<a:solidFill> 得到 rgb（支持 theme/prst）。"""
        try:
            el = sh._element
        except Exception:
            return None
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        ln = el.find(f".//{ns_a}ln")
        if ln is None:
            return None
        sf = ln.find(f"{ns_a}solidFill")
        return _resolve_fill_element_color(sf)

    def _shape_fill_color_rgb(sh):
        try:
            el = sh._element
        except Exception:
            return None
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        sp_pr = el.find(f".//{ns_a}spPr")
        if sp_pr is None:
            return None
        sf = sp_pr.find(f"{ns_a}solidFill")
        return _resolve_fill_element_color(sf)

    def _shape_has_no_fill(sh) -> bool:
        try:
            el = sh._element
        except Exception:
            return False
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        sp_pr = el.find(f".//{ns_a}spPr")
        if sp_pr is None:
            return True
        if sp_pr.find(f"{ns_a}noFill") is not None:
            return True
        return sp_pr.find(f"{ns_a}solidFill") is None

    def _shape_line_width(sh) -> float:
        """返回线宽（磅，EMU/12700）；无 <a:ln w="…"> 时返回 0。"""
        try:
            el = sh._element
        except Exception:
            return 0.0
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        ln = el.find(f".//{ns_a}ln")
        if ln is None:
            return 0.0
        w = ln.get("w")
        try:
            return int(w) / 12700.0 if w else 0.0
        except Exception:
            return 0.0

    def shape_color_is_light_blue(sh) -> bool:
        for rgb in (_shape_line_color_rgb(sh), _shape_fill_color_rgb(sh)):
            if is_light_blue(rgb):
                return True
        return False

    def rect_contains(outer, inner):
        ox, oy, ow, oh = outer
        ix, iy, iw, ih = inner
        return (ox <= ix and oy <= iy
                and ox + ow >= ix + iw and oy + oh >= iy + ih)

    def _rect_overlap_area(a, b) -> float:
        ax, ay, aw, ah = a
        bx, by, bw, bh = b
        ix = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
        iy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
        return ix * iy

    def _iter_editable_children(shape):
        """把 shape 展开为可编辑基础形状序列：
           GroupShape -> 递归其子对象；AUTO_SHAPE/FREEFORM 直接返回。
        """
        try:
            st = shape.shape_type
        except Exception:
            return []
        if st == MSO_SHAPE_TYPE.GROUP:
            out = []
            try:
                for sub in shape.shapes:
                    out.extend(_iter_editable_children(sub))
            except Exception:
                pass
            return out
        if st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
            return [shape]
        return []

    def _is_editable_border_segment(sh, title_rect) -> bool:
        """判定 sh 是"外围线框"而非填充覆盖：
           - 必须有线条颜色或线宽 > 0；
           - 若其矩形包围标题，必须是无填充/透明填充（避免覆盖标题）。
        """
        has_line = _shape_line_color_rgb(sh) is not None or _shape_line_width(sh) > 0
        if not has_line:
            return False
        if rect_contains(rect(sh), title_rect) and not _shape_has_no_fill(sh):
            # 填充色包围标题会遮挡标题
            return False
        return True

    title_keywords_by_slide = [
        ["高效动力链方案"],   # 第1页主标题
        ["整车性能协同优化"], # 第2页主标题
    ]

    page_pass = []
    page_detail = []
    for si, slide in enumerate(list(prs.slides)[:len(title_keywords_by_slide)]):
        # ① 找主标题文本框
        title_shape = None
        for idx, sh in enumerate(slide.shapes):
            if sh.has_text_frame and any(k in sh.text_frame.text for k in title_keywords_by_slide[si]):
                if sh.text_frame.text.strip():
                    title_shape = sh
                    break
        cond1 = title_shape is not None

        cond_border = False       # ② 有组合边框且包围
        cond_editable = False     # ③ 由 ≥2 段可编辑形状构成"外围线框"
        cond_color = False        # ④ 浅蓝色

        if cond1:
            t_rect = rect(title_shape)
            # 收集候选边框构件：
            #   路径 A：GroupShape，内含 ≥2 个 AUTO_SHAPE/FREEFORM 子对象；
            #   路径 B：多个独立 AUTO_SHAPE/FREEFORM 共同拼出边框。
            group_candidates = []
            for sh in slide.shapes:
                if sh is title_shape:
                    continue
                try:
                    st = sh.shape_type
                except Exception:
                    continue
                if st == MSO_SHAPE_TYPE.GROUP:
                    children = _iter_editable_children(sh)
                    if len(children) >= 2 and rect_contains(rect(sh), t_rect):
                        group_candidates.append((sh, children))

            border_segments = []
            for sh in slide.shapes:
                if sh is title_shape:
                    continue
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    continue
                try:
                    st = sh.shape_type
                except Exception:
                    continue
                if st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
                    # 相邻 / 包围标题的线框构件
                    if _rect_overlap_area(rect(sh), t_rect) > 0 or rect_contains(rect(sh), t_rect):
                        if _is_editable_border_segment(sh, t_rect):
                            border_segments.append(sh)

            # 判定：GroupShape 优先
            if group_candidates:
                grp, kids = group_candidates[0]
                kids_ok = [k for k in kids if _is_editable_border_segment(k, t_rect)]
                if len(kids_ok) >= 2:
                    cond_border = True
                    cond_editable = True
                    if any(shape_color_is_light_blue(k) for k in kids_ok):
                        cond_color = True

            # 若没有 GroupShape，尝试用多个独立可编辑形状构造组合边框
            if not cond_border and len(border_segments) >= 2:
                # 组合外接矩形：取所有 border_segments 的并集矩形
                min_left = min(cm(s.left) for s in border_segments)
                min_top = min(cm(s.top) for s in border_segments)
                max_right = max(cm(s.left) + cm(s.width) for s in border_segments)
                max_bottom = max(cm(s.top) + cm(s.height) for s in border_segments)
                combined_rect = (min_left, min_top,
                                 max_right - min_left, max_bottom - min_top)
                if rect_contains(combined_rect, t_rect):
                    cond_border = True
                    cond_editable = True
                    if any(shape_color_is_light_blue(s) for s in border_segments):
                        cond_color = True

        ok_page = cond1 and cond_border and cond_editable and cond_color
        page_pass.append(ok_page)
        page_detail.append(
            f"slide{si+1}[标题存在={cond1},组合边框包围={cond_border},"
            + f"可编辑构件≥2={cond_editable},浅蓝色={cond_color}]"
        )

    hit_border = all(page_pass)
    points.append((+3,
                   "第1、2页顶部主标题：外围添加浅蓝色组合边框，边框为可编辑形状",
                   "顶部主标题外围浅蓝色可编辑组合边框：" + " | ".join(page_detail),
                   hit_border))

    return points


def _find_target_file(dir_path: str) -> str:
    """在脚本所在目录里定位被评估的 .pptx 文件（仅识别 .pptx，不再解析 .ppt）。"""
    if not os.path.isdir(dir_path):
        return ""
    candidates = []
    for name in os.listdir(dir_path):
        low = name.lower()
        if low.endswith(".pptx"):
            candidates.append(os.path.join(dir_path, name))
    if not candidates:
        return ""
    # 优先匹配模板文件名
    preferred = os.path.join(dir_path, "核心技术_照片排版版.pptx")
    if preferred in candidates:
        return preferred
    return candidates[0]


def evaluate(dir_path: str):
    """统一入口：接收脚本所在目录路径，脚本自行定位待评估文档并返回结构化结果。"""
    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": 0,
    }
    try:
        file_path = _find_target_file(dir_path)
        if not file_path:
            result["status"] = "error"
            result["error"] = f"目录中未找到 .pptx 文件：{dir_path}"
            return result
        result["file_name"] = os.path.basename(file_path)

        ok, reasons = check_dimension1(file_path)
        result["dim1_pass"] = ok
        if not ok:
            result["dim1_reason"] = "；".join(reasons)
            result["total_score"] = 0
            # 维度一未通过时不再展开维度二逐项
            return result

        points = check_dimension2(file_path)
        dim2_items = []
        max_score = 0
        total_score = 0
        for delta, rule, _detail, hit in points:
            # 正向得分项：max_delta = +delta，未命中 delta=0；
            # 扣分项：max_delta = 负值（表示最大扣分幅度），未命中不扣分（delta=0），命中则应用扣分（delta=max_delta）。
            max_delta = delta
            actual = delta if hit else 0
            if max_delta > 0:
                max_score += max_delta
            total_score += actual
            dim2_items.append({
                "rule": rule,
                "max_delta": max_delta,
                "delta": actual,
                "hit": bool(hit),
                "detail": "",
            })

        result["dim2_items"] = dim2_items
        result["total_score"] = total_score
        result["max_score"] = max_score
        return result
    except Exception as e:
        result["status"] = "error"
        result["error"] = f"{type(e).__name__}: {e}"
        return result


if __name__ == "__main__":
    # 本地自测：默认以脚本所在目录作为 dir_path，也支持命令行显式传入
    _dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(_dir), ensure_ascii=False, indent=2))
