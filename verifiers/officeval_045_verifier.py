"""
海岚清齿声波护理品牌汇报 PPT 自动评估脚本
评分逻辑：
1. 维度1（可用与可修改性）：不满足直接判0分
2. 维度2（完成度）：得分点 + 扣分点累计

对外接口：``evaluate(dir_path: str) -> dict``
- ``dir_path`` 为脚本所在目录，脚本自行在该目录下定位并打开被评估的 .pptx 文件
- 返回结构见《脚本接口差异与统一建议.md》§2.2
"""
import sys

SCRIPT_ID = "045"

import os
import json
import re
import zipfile

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree
from PIL import Image

EMU_PER_CM = 360000
SLIDE_WIDTH_CM = 33.86
SLIDE_HEIGHT_CM = 19.05
NS = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
      'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

# 由 evaluate() 在运行时填充；下方辅助函数以模块全局形式引用
prs = None
slides = None
slide_w = None
slide_h = None



def emu_to_cm(emu):
    return emu / EMU_PER_CM


def get_all_text(slide):
    """获取页面所有文本"""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    return ' '.join(texts)


def get_text_shapes(slide):
    """获取所有文本shape"""
    return [s for s in slide.shapes if s.has_text_frame]


def get_media_shapes(slide):
    """获取所有媒体shape"""
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.MEDIA]


def get_picture_shapes(slide):
    """获取所有图片shape"""
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


def shape_overlaps(s1, s2):
    """检查两个shape是否重叠"""
    l1, t1, r1, b1 = s1.left, s1.top, s1.left + s1.width, s1.top + s1.height
    l2, t2, r2, b2 = s2.left, s2.top, s2.left + s2.width, s2.top + s2.height
    return not (r1 <= l2 or r2 <= l1 or b1 <= t2 or b2 <= t1)


def shape_out_of_bounds(shape, slide_w, slide_h):
    """检查shape是否超出页面"""
    r = shape.left + shape.width
    b = shape.top + shape.height
    margin = Cm(0.5)  # 允许小误差
    return shape.left < -margin or shape.top < -margin or r > slide_w + margin or b > slide_h + margin


def shape_covers_slide(shape, tolerance=Cm(0.2)):
    """检查shape是否覆盖完整页面。"""
    return (shape.left <= tolerance and shape.top <= tolerance and
            shape.left + shape.width >= slide_w - tolerance and
            shape.top + shape.height >= slide_h - tolerance)


def extract_color_rgb(color_parent):
    """从Office OOXML颜色节点读取RGB值。"""
    srgb = color_parent.find('./a:srgbClr', NS)
    if srgb is not None and srgb.get('val'):
        val = srgb.get('val')
        if re.fullmatch(r'[0-9A-Fa-f]{6}', val):
            return tuple(int(val[i:i + 2], 16) for i in (0, 2, 4))

    sys_clr = color_parent.find('./a:sysClr', NS)
    if sys_clr is not None and sys_clr.get('lastClr'):
        val = sys_clr.get('lastClr')
        if re.fullmatch(r'[0-9A-Fa-f]{6}', val):
            return tuple(int(val[i:i + 2], 16) for i in (0, 2, 4))

    scrgb = color_parent.find('./a:scrgbClr', NS)
    if scrgb is not None:
        try:
            return tuple(round(int(scrgb.get(c, '0')) * 255 / 100000) for c in ('r', 'g', 'b'))
        except ValueError:
            return None
    return None


def is_white(rgb):
    return rgb is not None and rgb[0] >= 245 and rgb[1] >= 245 and rgb[2] >= 245


def is_light_blue(rgb):
    if rgb is None:
        return False
    r, g, b = rgb
    return r >= 170 and g >= 210 and b >= 230 and b >= g and g >= r - 20


def gradfill_is_linear_lightblue_to_white(grad_fill):
    """检查填充是否为浅蓝到白色的线性渐变。"""
    if grad_fill is None or grad_fill.find('./a:lin', NS) is None:
        return False
    stops = []
    for idx, stop in enumerate(grad_fill.findall('./a:gsLst/a:gs', NS)):
        rgb = extract_color_rgb(stop)
        if rgb is not None:
            try:
                pos = int(stop.get('pos', idx))
            except ValueError:
                pos = idx
            stops.append((pos, rgb))
    stops.sort(key=lambda item: item[0])
    first_light_blue = next((i for i, (_, rgb) in enumerate(stops) if is_light_blue(rgb)), None)
    last_white = next((i for i in range(len(stops) - 1, -1, -1) if is_white(stops[i][1])), None)
    return first_light_blue is not None and last_white is not None and first_light_blue < last_white


def slide_has_office_linear_lightblue_white_bg(slide):
    """检查第16页背景是否为Office可识别的整页浅蓝到白色线性渐变。"""
    bg_pr = slide._element.find('./p:cSld/p:bg/p:bgPr', NS)
    if bg_pr is not None:
        grad_fill = bg_pr.find('./a:gradFill', NS)
        if gradfill_is_linear_lightblue_to_white(grad_fill):
            return True

    for shape in slide.shapes:
        if not shape_covers_slide(shape):
            continue
        grad_fill = shape._element.find('./p:spPr/a:gradFill', NS)
        if gradfill_is_linear_lightblue_to_white(grad_fill):
            return True
    return False


def _slide_has_picture_watermark(slide) -> bool:
    """判断第16页是否用"图片水印"替代渐变背景。
    命中条件：
      - slide background 使用图片填充（<p:bg><p:bgPr><a:blipFill/>）；或
      - 页面里存在覆盖整页的图片形状（<p:pic> 且尺寸接近整页）。
    仅整页覆盖的图片视作"替代渐变"的水印；页面上的装饰性小图片不算。
    """
    bg_pr = slide._element.find('./p:cSld/p:bg/p:bgPr', NS)
    if bg_pr is not None and bg_pr.find('./a:blipFill', NS) is not None:
        return True
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape_covers_slide(shape):
            return True
    return False


def _slide_has_dark_noise_or_solid_bg(slide) -> bool:
    """判断第16页是否用"深色杂色 / 纯色背景"替代渐变背景。
    命中条件（任一）：
      - slide background 使用 <a:solidFill>，颜色非浅蓝且非白色（视为纯色替代）。
      - slide background 使用 <a:pattFill>（图案填充，OOXML 里"杂色/网点"背景
        的典型写法）。
      - 存在整页覆盖形状使用 <a:solidFill> 且颜色非浅蓝非白色，或使用
        <a:pattFill>。
    浅蓝→白线性渐变（<a:gradFill>）不在本函数命中范围。
    """
    bg_pr = slide._element.find('./p:cSld/p:bg/p:bgPr', NS)
    if bg_pr is not None:
        if bg_pr.find('./a:pattFill', NS) is not None:
            return True
        solid = bg_pr.find('./a:solidFill', NS)
        if solid is not None:
            rgb = extract_color_rgb(solid)
            if rgb is not None and not (is_light_blue(rgb) or is_white(rgb)):
                return True

    for shape in slide.shapes:
        if not shape_covers_slide(shape):
            continue
        sp_pr = shape._element.find('./p:spPr', NS)
        if sp_pr is None:
            continue
        if sp_pr.find('./a:pattFill', NS) is not None:
            return True
        solid = sp_pr.find('./a:solidFill', NS)
        if solid is not None:
            rgb = extract_color_rgb(solid)
            if rgb is not None and not (is_light_blue(rgb) or is_white(rgb)):
                return True
    return False


def get_font_info(shape):
    """获取shape中的字体信息"""
    fonts = []
    sizes = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.name:
                    fonts.append(run.font.name)
                if run.font.size:
                    sizes.append(run.font.size)
    return fonts, sizes


def is_sans_serif(font_name):
    """判断是否为无衬线字体"""
    if not font_name:
        return True
    sans_fonts = ['微软雅黑', 'Microsoft YaHei', 'Arial', 'Helvetica',
                  'Calibri', 'Segoe UI', 'SimHei', '黑体', 'Noto Sans',
                  'PingFang SC', 'Source Han Sans', 'DengXian', '等线']
    return any(f.lower() in font_name.lower() or font_name.lower() in f.lower()
               for f in sans_fonts)


def evaluate(dir_path: str) -> dict:
    """在 ``dir_path`` 目录内定位 .pptx 文件并进行评估，返回结构化结果字典。

    参数:
        dir_path: 脚本所在目录路径。脚本会在此目录内自行查找并打开被评估文档。

    返回:
        参见《脚本接口差异与统一建议.md》§2.2 定义的结构化字典。
    """
    global prs, slides, slide_w, slide_h  # 供顶层辅助函数以模块全局形式引用

    script_id = "045"
    # 在 dir_path 内查找首个 .pptx 文件（只识别 .pptx，不再兼容老版 .ppt 二进制格式；
    # 忽略 Office 临时文件 ~$*）
    pptx_path = None
    file_name = ""
    if os.path.isdir(dir_path):
        for name in sorted(os.listdir(dir_path)):
            if name.startswith("~$"):
                continue
            if name.lower().endswith(".pptx"):
                pptx_path = os.path.join(dir_path, name)
                file_name = name
                break
    if pptx_path is None:
        return {
            "id": script_id,
            "file_name": "",
            "status": "error",
            "error": f"在目录 {dir_path} 下未找到 .pptx 文件",
            "dim1_pass": False,
            "dim1_reason": "未找到被评估文档",
            "dim2_items": [],
            "total_score": 0,
            "max_score": 0,
        }

    results = []  # (points, description, triggered, msg)
    # ========== 维度1：可用与可修改性 ==========
    dim1_pass = True
    dim1_issues = []

    # 1.1 文件格式和可打开
    if not pptx_path.lower().endswith('.pptx'):
        dim1_pass = False
        dim1_issues.append("文件不是.pptx格式")

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        dim1_pass = False
        dim1_issues.append(f"文件无法正常打开: {e}")
        return {
            "id": script_id,
            "file_name": file_name,
            "status": "error",
            "error": '; '.join(dim1_issues),
            "dim1_pass": False,
            "dim1_reason": '; '.join(dim1_issues),
            "dim2_items": [],
            "total_score": 0,
            "max_score": 0,
        }

    # 1.2 页数为17页
    slide_count = len(prs.slides)
    if slide_count != 17:
        dim1_pass = False
        dim1_issues.append(f"页数为{slide_count}页，非17页")

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    if not dim1_pass:
        return {
            "id": script_id,
            "file_name": file_name,
            "status": "ok",
            "error": None,
            "dim1_pass": False,
            "dim1_reason": '; '.join(dim1_issues),
            "dim2_items": [],
            "total_score": 0,
            "max_score": 0,
        }

    # ========== 维度2：完成度评分 ==========
    slides = list(prs.slides)

    # --- +3: 第16页新增视频页 ---
    def check_slide16_new_video_page():
        """第16页新增视频页，位于原结束页之前，第17页仍为结束页；
        第16页整页背景为线性渐变浅蓝到白色；无深色杂色/水印/纯色替代"""
        slide16 = slides[15]
        slide17 = slides[16]
        # 第17页应为结束页（含"感谢观看"）
        s17_text = get_all_text(slide17)
        if '感谢观看' not in s17_text:
            return False, "第17页不是结束页（无'感谢观看'）"
        # 第16页有视频
        media = get_media_shapes(slide16)
        if len(media) == 0:
            return False, "第16页无视频对象"
        # 第16页背景：明确禁止的替代实现
        # (1) 图片水印替代（<a:blipFill> 背景，或整页覆盖的图片形状）
        if _slide_has_picture_watermark(slide16):
            return False, "第16页背景使用图片水印替代渐变"
        # (2) 深色杂色 / 纯色背景替代（<a:pattFill>；或非浅蓝非白色的 <a:solidFill>）
        if _slide_has_dark_noise_or_solid_bg(slide16):
            return False, "第16页背景存在深色杂色/纯色替代渐变"
        # (3) 允许的渐变实现：Office 可识别的整页 <a:gradFill>，
        #     从浅蓝色过渡到白色的线性渐变（<a:lin>），色标顺序为浅蓝在前、白色在后。
        if not slide_has_office_linear_lightblue_white_bg(slide16):
            return False, "第16页背景不是覆盖整页的浅蓝到白色线性渐变"
        return True, ""

    passed, msg = check_slide16_new_video_page()
    results.append((3, "第16页新增视频页+渐变背景+结束页位置", passed, msg))

    # --- +5: 第16页三个视频整体排列 ---
    def check_slide16_video_layout():
        """三个视频从左到右依次为清洁、定时、护龈，横向居中等距排列，
        视频内容对应；视频为可播放视频对象，不是静态图片；宽高比接近1:2"""
        slide16 = slides[15]
        media = get_media_shapes(slide16)

        if len(media) != 3:
            return False, f"第16页视频数量为{len(media)}，非3个"

        ns_r = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        VIDEO_REL_TYPES = {
            f'{ns_r}/video',
            'http://schemas.microsoft.com/office/2007/relationships/media',
        }
        VIDEO_EXTS = re.compile(r'\.(mp4|mov|m4v|wmv|avi|mpeg|mpg)$', re.I)

        def get_ooxml_name(shape):
            c_nv_pr = shape._element.find('.//p:cNvPr', NS)
            return c_nv_pr.get('name', '') if c_nv_pr is not None else shape.name

        def is_playable_video(shape):
            """通过OOXML a:videoFile + 关系类型确认是可播放视频对象，而非静态图片。"""
            video_node = shape._element.find('.//p:nvPr/a:videoFile', NS)
            if video_node is None:
                return False
            rel_id = video_node.get(f'{{{ns_r}}}link')
            if not rel_id:
                return False
            rel = shape.part.rels.get(rel_id)
            if rel is None:
                return False
            if rel.reltype not in VIDEO_REL_TYPES and not rel.reltype.endswith('/video'):
                return False
            target = getattr(rel, 'target_ref', '') or ''
            return bool(VIDEO_EXTS.search(target))

        def content_keyword(shape):
            """从cNvPr/@name中提取视频内容标识；Office将原始文件名存在此处。"""
            name = get_ooxml_name(shape).lower()
            name_stripped = re.sub(r'[\s_\-.()\d]+', '', name)
            if any(k in name_stripped for k in ('clean', 'cleaning', '清洁')):
                return '清洁'
            if any(k in name_stripped for k in ('timer', 'timing', '定时')):
                return '定时'
            if any(k in name_stripped for k in ('gumcare', 'gum', '护龈')):
                return '护龈'
            return None

        # 1. 确认每个媒体对象都是可播放视频，而非静态图片
        for m in media:
            if not is_playable_video(m):
                return False, f"'{get_ooxml_name(m)}'不含a:videoFile关系，不是可播放视频对象"

        # 2. 按水平位置从左到右排序，验证顺序为清洁→定时→护龈（视频内容对应）
        media_sorted = sorted(media, key=lambda s: s.left)
        expected_order = ['清洁', '定时', '护龈']
        for m, expected in zip(media_sorted, expected_order):
            actual = content_keyword(m)
            if actual != expected:
                return False, (f"视频从左到右内容顺序不正确，应为清洁、定时、护龈；"
                               f"实际={[content_keyword(s) or get_ooxml_name(s) for s in media_sorted]}")

        # 3. 宽高比接近1:2（宽:高≈0.5，允许范围0.3-0.7）
        for m in media_sorted:
            ratio = m.width / m.height if m.height > 0 else 0
            if not (0.3 <= ratio <= 0.7):
                return False, f"视频'{get_ooxml_name(m)}'宽高比{ratio:.2f}不接近1:2"

        # 4. 等距排列：相邻视频间隔差不超过0.2cm
        gap1 = media_sorted[1].left - (media_sorted[0].left + media_sorted[0].width)
        gap2 = media_sorted[2].left - (media_sorted[1].left + media_sorted[1].width)
        if gap1 <= 0 or gap2 <= 0:
            return False, "视频存在重叠，无法构成等距排列"
        if abs(emu_to_cm(gap1) - emu_to_cm(gap2)) > 0.2:
            return False, f"视频间距不等: gap1={emu_to_cm(gap1):.2f}cm, gap2={emu_to_cm(gap2):.2f}cm"

        # 5. 横向居中：三视频整体中心与页面中心偏差不超过0.3cm
        group_left = media_sorted[0].left
        group_right = media_sorted[-1].left + media_sorted[-1].width
        group_center = (group_left + group_right) / 2
        if abs(group_center - slide_w / 2) > Cm(0.3):
            return False, (f"三视频整体未横向居中: "
                           f"组中心={emu_to_cm(group_center):.2f}cm, "
                           f"页面中心={emu_to_cm(slide_w/2):.2f}cm")

        return True, ""

    passed, msg = check_slide16_video_layout()
    results.append((5, "第16页三个视频整体排列(清洁/定时/护龈+等距+宽高比1:2)", passed, msg))

    # --- +5/+5/+5: 第16页视频动画（点击触发顺序）---
    def check_slide16_animations():
        """检查视频点击触发动画；左侧项要求第一下点击出现并播放清洁视频。
        返回 (三项结果, 三项独立消息)：每项消息仅反映该项自身报错/跳过的理由，
        通过时对应位置为空串，避免不同项之间的失败原因互相污染。"""
        slide16 = slides[15]
        timing = slide16._element.find('.//p:timing', NS)
        if timing is None:
            return [False, False, False], ["无动画"] * 3

        media = sorted(get_media_shapes(slide16), key=lambda s: s.left)
        if len(media) < 3:
            return [False, False, False], ["第16页视频不足3个"] * 3

        def shape_id(shape):
            c_nv_pr = shape._element.find('.//p:cNvPr', NS)
            return c_nv_pr.get('id') if c_nv_pr is not None else str(shape.shape_id)

        def shape_name(shape):
            c_nv_pr = shape._element.find('.//p:cNvPr', NS)
            return (c_nv_pr.get('name', '') if c_nv_pr is not None else shape.name).lower()

        def video_keyword(shape):
            name = re.sub(r'[\s_\-.()\d]+', '', shape_name(shape))
            if 'clean' in name or 'cleaning' in name or '清洁' in name:
                return '清洁'
            if 'timer' in name or 'timing' in name or '定时' in name:
                return '定时'
            if 'gumcare' in name or 'gum' in name or '护龈' in name:
                return '护龈'
            return None

        def is_clean_video(shape):
            return video_keyword(shape) == '清洁'

        def is_timer_video(shape):
            return video_keyword(shape) == '定时'

        def is_gum_video(shape):
            return video_keyword(shape) == '护龈'

        def rect(shape):
            return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)

        def rect_center_inside(inner, outer):
            i = rect(inner)
            o = rect(outer)
            cx = (i[0] + i[2]) / 2
            cy = (i[1] + i[3]) / 2
            return o[0] <= cx <= o[2] and o[1] <= cy <= o[3]

        def related_card_spids(video_shape):
            ids = {shape_id(video_shape)}
            for shape in slide16.shapes:
                if shape is video_shape:
                    continue
                if rect_center_inside(shape, video_shape) or rect_center_inside(video_shape, shape):
                    sid = shape_id(shape)
                    if sid:
                        ids.add(sid)
            return ids

        left_video = media[0]
        if not is_clean_video(left_video):
            return [False, False, False], ["左侧视频不是清洁视频"] * 3

        # 找mainSeq中的par序列；每个顶层par对应一次点击步骤
        main_seq = timing.find('.//p:seq/p:cTn[@nodeType="mainSeq"]', NS)
        if main_seq is None:
            return [False, False, False], ["无主序列动画"] * 3

        child_pars = main_seq.findall('./p:childTnLst/p:par', NS)
        if not child_pars:
            return [False, False, False], ["无点击动画步骤"] * 3

        def par_is_click_triggered(par):
            ctn = par.find('./p:cTn', NS)
            if ctn is None:
                return False
            cond = ctn.find('./p:stCondLst/p:cond', NS)
            return cond is not None and cond.get('delay') == 'indefinite'

        def par_targets(par):
            return [t.get('spid') for t in par.findall('.//p:spTgt', NS) if t.get('spid')]

        def par_has_appear_for(par, allowed_spids):
            for anim in par.findall('.//p:animEffect', NS):
                if anim.get('transition') != 'in':
                    continue
                targets = [t.get('spid') for t in anim.findall('.//p:spTgt', NS)]
                if any(spid in allowed_spids for spid in targets):
                    return True
            return False

        def par_has_video_play_for(par, video_spid):
            for video in par.findall('.//p:video', NS):
                targets = [t.get('spid') for t in video.findall('.//p:spTgt', NS)]
                if video_spid in targets:
                    return True
            return False

        def par_has_any_target_for(par, allowed_spids):
            targets = par_targets(par)
            return any(spid in allowed_spids for spid in targets)

        def click_step_ok(par, video_shape, allowed_spids):
            return (par_is_click_triggered(par) and
                    par_has_appear_for(par, allowed_spids) and
                    par_has_video_play_for(par, shape_id(video_shape)))

        left_allowed_spids = related_card_spids(left_video)
        first_click = child_pars[0]
        left_ok = bool(click_step_ok(first_click, left_video, left_allowed_spids))

        mid_video = media[1]
        right_video = media[2]
        mid_spid = shape_id(mid_video)
        right_spid = shape_id(right_video)
        mid_allowed_spids = related_card_spids(mid_video)
        right_allowed_spids = related_card_spids(right_video)

        # 中间项细则：点击第二下出现并播放；顺序位于左侧清洁视频之后、右侧护龈视频之前；
        # 播放对象必须是中间timer视频，出现对象可为该视频或其对应卡片。
        second_click = child_pars[1] if len(child_pars) >= 2 else None
        third_click = child_pars[2] if len(child_pars) >= 3 else None
        mid_ok = bool(second_click is not None and third_click is not None and
                      is_timer_video(mid_video) and is_gum_video(right_video) and
                      click_step_ok(second_click, mid_video, mid_allowed_spids) and
                      par_has_any_target_for(first_click, left_allowed_spids) and
                      par_has_any_target_for(second_click, mid_allowed_spids) and
                      par_has_any_target_for(third_click, right_allowed_spids) and
                      par_has_video_play_for(third_click, right_spid))

        # 右侧项细则：点击第三下出现并播放；触发顺序位于左侧清洁视频和中间定时视频之后；
        # 播放对象必须是右侧gumcare视频，出现对象可为该视频或其对应卡片。
        right_ok = bool(second_click is not None and third_click is not None and
                        is_clean_video(left_video) and is_timer_video(mid_video) and is_gum_video(right_video) and
                        par_has_any_target_for(first_click, left_allowed_spids) and
                        par_has_any_target_for(second_click, mid_allowed_spids) and
                        click_step_ok(third_click, right_video, right_allowed_spids))

        left_msg = "" if left_ok else "左侧清洁视频未在第一下点击后出现并播放"
        mid_msg = "" if mid_ok else "中间timer视频未在第二下点击后出现并播放（顺序应在清洁之后、护龈之前）"
        right_msg = "" if right_ok else "右侧gumcare视频未在第三下点击后出现并播放（顺序应在清洁和定时之后）"

        return [left_ok, mid_ok, right_ok], [left_msg, mid_msg, right_msg]

    anim_results, anim_msgs = check_slide16_animations()
    results.append((5, "第16页左侧视频动画(点击第一下出现播放清洁视频)", anim_results[0], "" if anim_results[0] else anim_msgs[0]))
    results.append((5, "第16页中间视频动画(点击第二下出现播放定时视频)", anim_results[1], "" if anim_results[1] else anim_msgs[1]))
    results.append((5, "第16页右侧视频动画(点击第三下出现播放护龈视频)", anim_results[2], "" if anim_results[2] else anim_msgs[2]))

    # --- +1: 第9页主体图片 ---
    def check_slide9_main_image():
        """第9页主体图片：位于页面中部，尺寸宽17.60cm/高16.82cm；
        等比例缩放；边缘不超出页面可视区域；
        不遮挡标题'03 产品生命周期管理'和页脚品牌文字'海岚清齿 HAILAN CARE'。"""
        slide = slides[8]
        pics = get_picture_shapes(slide)
        EXPECTED_W = 17.60
        EXPECTED_H = 16.82
        EXPECTED_RATIO = EXPECTED_W / EXPECTED_H  # 原始比例 ≈ 1.046
        RATIO_TOLERANCE = 0.05   # 等比例允许误差
        SIZE_TOLERANCE = 0.5     # cm，尺寸匹配允许误差

        target = None
        for p in pics:
            w = emu_to_cm(p.width)
            h = emu_to_cm(p.height)
            if abs(w - EXPECTED_W) <= SIZE_TOLERANCE and abs(h - EXPECTED_H) <= SIZE_TOLERANCE:
                target = p
                break
        if target is None:
            return False, f"未找到尺寸约为{EXPECTED_W}x{EXPECTED_H}cm的主体图片"

        w = emu_to_cm(target.width)
        h = emu_to_cm(target.height)

        # 1. 等比例缩放：实际宽高比与原始比例的偏差在容差内
        actual_ratio = w / h if h > 0 else 0
        if abs(actual_ratio - EXPECTED_RATIO) > RATIO_TOLERANCE:
            return False, f"图片宽高比{actual_ratio:.3f}与等比例{EXPECTED_RATIO:.3f}偏差过大，未等比缩放"

        # 2. 边缘不超出页面可视区域（四边都检查）
        margin = Cm(0.3)
        if target.left < -margin:
            return False, "图片左边缘超出页面"
        if target.top < -margin:
            return False, "图片上边缘超出页面"
        if target.left + target.width > slide_w + margin:
            return False, "图片右边缘超出页面"
        if target.top + target.height > slide_h + margin:
            return False, "图片下边缘超出页面"

        # 3. 位于页面中部：主体图片中心位于页面可视区域的中间带，避免贴边/页眉/页脚；
        #    不强制水平居中，因为该页左侧为标题说明区，主体图片按细则只要求位于页面中部。
        img_cx = emu_to_cm(target.left + target.width / 2)
        img_cy = emu_to_cm(target.top + target.height / 2)
        if not (SLIDE_WIDTH_CM * 0.35 <= img_cx <= SLIDE_WIDTH_CM * 0.95):
            return False, f"图片水平中心{img_cx:.2f}cm不在页面中部区域"
        if not (SLIDE_HEIGHT_CM * 0.20 <= img_cy <= SLIDE_HEIGHT_CM * 0.80):
            return False, f"图片垂直中心{img_cy:.2f}cm不在页面中部区域"

        # 4. 不遮挡标题"03 产品生命周期管理"和页脚品牌文字"海岚清齿 HAILAN CARE"
        #    遮挡判定：shape_overlaps为真，且图片z-order高于文本（覆盖在上方）
        shapes_list = list(slide.shapes)
        p_z = shapes_list.index(target) if target in shapes_list else -1
        PROTECTED_KEYWORDS = ["03 产品生命周期管理", "海岚清齿", "HAILAN CARE"]
        for s in get_text_shapes(slide):
            text = s.text_frame.text
            if any(kw in text for kw in PROTECTED_KEYWORDS):
                if shape_overlaps(target, s):
                    s_z = shapes_list.index(s) if s in shapes_list else -1
                    if p_z > s_z:  # 图片在文本上层 → 遮挡
                        excerpt = text.strip()[:20]
                        return False, f"图片遮挡了文本'{excerpt}'"

        return True, ""

    passed9, msg9 = check_slide9_main_image()
    results.append((1, "第9页主体图片(17.60x16.82cm+中部+等比+不超出+不遮挡标题页脚)", passed9, msg9))


    def check_slide12_main_image():
        """第12页主体图片：位于页面中部，尺寸宽17.60cm/高16.82cm；
        图片保持等比例缩放；边缘不超出页面可视区域。"""
        slide = slides[11]
        pics = get_picture_shapes(slide)
        EXPECTED_W = 17.60
        EXPECTED_H = 16.82
        EXPECTED_RATIO = EXPECTED_W / EXPECTED_H  # ≈ 1.046
        RATIO_TOLERANCE = 0.05
        SIZE_TOLERANCE = 0.5  # cm

        target = None
        for p in pics:
            w = emu_to_cm(p.width)
            h = emu_to_cm(p.height)
            if abs(w - EXPECTED_W) <= SIZE_TOLERANCE and abs(h - EXPECTED_H) <= SIZE_TOLERANCE:
                target = p
                break
        if target is None:
            return False, f"未找到尺寸约为{EXPECTED_W}x{EXPECTED_H}cm的主体图片"

        w = emu_to_cm(target.width)
        h = emu_to_cm(target.height)

        # 1. 等比例缩放：实际宽高比与原始比例偏差在容差内
        actual_ratio = w / h if h > 0 else 0
        if abs(actual_ratio - EXPECTED_RATIO) > RATIO_TOLERANCE:
            return False, f"图片宽高比{actual_ratio:.3f}与等比例{EXPECTED_RATIO:.3f}偏差过大，未等比缩放"

        # 2. 边缘不超出页面可视区域（四边均检查）
        margin = Cm(0.3)
        if target.left < -margin:
            return False, "图片左边缘超出页面"
        if target.top < -margin:
            return False, "图片上边缘超出页面"
        if target.left + target.width > slide_w + margin:
            return False, "图片右边缘超出页面"
        if target.top + target.height > slide_h + margin:
            return False, "图片下边缘超出页面"

        # 3. 位于页面中部：图片中心在页面宽35%以内或以右均可（左图右文布局），
        #    垂直中心在页面高20%~85%之间
        img_cx = emu_to_cm(target.left + target.width / 2)
        img_cy = emu_to_cm(target.top + target.height / 2)
        if not (SLIDE_WIDTH_CM * 0.05 <= img_cx <= SLIDE_WIDTH_CM * 0.75):
            return False, f"图片水平中心{img_cx:.2f}cm不在页面中部区域"
        if not (SLIDE_HEIGHT_CM * 0.20 <= img_cy <= SLIDE_HEIGHT_CM * 0.85):
            return False, f"图片垂直中心{img_cy:.2f}cm不在页面中部区域"

        return True, ""

    passed12, msg12 = check_slide12_main_image()
    results.append((1, "第12页主体图片(17.60x16.82cm+中部+等比+不超出)", passed12, msg12))


    def check_slide14_main_image():
        """第14页主体图片：位于页面中部，尺寸宽17.60cm/高16.82cm；
        等比例缩放；边缘不超出页面可视区域；
        不遮挡'05 产品组合'、'ORAL CARE BRAND'和页脚品牌文字'海岚清齿 HAILAN CARE'。"""
        slide = slides[13]
        pics = get_picture_shapes(slide)
        EXPECTED_W = 17.60
        EXPECTED_H = 16.82
        EXPECTED_RATIO = EXPECTED_W / EXPECTED_H
        RATIO_TOLERANCE = 0.05
        SIZE_TOLERANCE = 0.5

        target = None
        for p in pics:
            w = emu_to_cm(p.width)
            h = emu_to_cm(p.height)
            if abs(w - EXPECTED_W) <= SIZE_TOLERANCE and abs(h - EXPECTED_H) <= SIZE_TOLERANCE:
                target = p
                break
        if target is None:
            return False, f"未找到尺寸约为{EXPECTED_W}x{EXPECTED_H}cm的主体图片"

        w = emu_to_cm(target.width)
        h = emu_to_cm(target.height)

        # 1. 等比例缩放：实际宽高比与原始比例偏差在容差内
        actual_ratio = w / h if h > 0 else 0
        if abs(actual_ratio - EXPECTED_RATIO) > RATIO_TOLERANCE:
            return False, f"图片宽高比{actual_ratio:.3f}与等比例{EXPECTED_RATIO:.3f}偏差过大，未等比缩放"

        # 2. 边缘不超出页面可视区域（四边均检查）
        margin = Cm(0.3)
        if target.left < -margin:
            return False, "图片左边缘超出页面"
        if target.top < -margin:
            return False, "图片上边缘超出页面"
        if target.left + target.width > slide_w + margin:
            return False, "图片右边缘超出页面"
        if target.top + target.height > slide_h + margin:
            return False, "图片下边缘超出页面"

        # 3. 位于页面中部：主体图片中心位于页面可视区域的中间带，适配左文右图版式。
        img_cx = emu_to_cm(target.left + target.width / 2)
        img_cy = emu_to_cm(target.top + target.height / 2)
        if not (SLIDE_WIDTH_CM * 0.35 <= img_cx <= SLIDE_WIDTH_CM * 0.95):
            return False, f"图片水平中心{img_cx:.2f}cm不在页面中部区域"
        if not (SLIDE_HEIGHT_CM * 0.20 <= img_cy <= SLIDE_HEIGHT_CM * 0.80):
            return False, f"图片垂直中心{img_cy:.2f}cm不在页面中部区域"

        # 4. 不遮挡指定文本。Office里以后层对象覆盖前层对象，且存在显著重叠，才判定为遮挡。
        def overlap_ratio(s1, s2):
            left = max(s1.left, s2.left)
            top = max(s1.top, s2.top)
            right = min(s1.left + s1.width, s2.left + s2.width)
            bottom = min(s1.top + s1.height, s2.top + s2.height)
            overlap = max(0, right - left) * max(0, bottom - top)
            text_area = s2.width * s2.height
            return overlap / text_area if text_area > 0 else 0

        shapes_list = list(slide.shapes)
        p_z = shapes_list.index(target) if target in shapes_list else -1
        protected_keywords = ["05", "产品组合", "ORAL CARE BRAND", "海岚清齿", "HAILAN CARE"]
        for s in get_text_shapes(slide):
            text = s.text_frame.text
            if any(kw in text for kw in protected_keywords):
                if shape_overlaps(target, s):
                    s_z = shapes_list.index(s) if s in shapes_list else -1
                    if p_z > s_z and overlap_ratio(target, s) > 0.05:
                        excerpt = text.strip()[:20]
                        return False, f"图片遮挡了文本'{excerpt}'"

        return True, ""

    passed14, msg14 = check_slide14_main_image()
    results.append((1, "第14页主体图片(17.60x16.82cm+中部+等比+不超出+不遮挡指定文字)", passed14, msg14))

    # --- +5: 页码文本 ---
    def check_page_numbers():
        """第2,4,5,7,8,10,12,13,15,16页右下角需有页码文本，且不遮挡页面内容。"""
        need_pages = [2, 4, 5, 7, 8, 10, 12, 13, 15, 16]  # 1-indexed

        def find_page_number_shape(slide):
            """在页面右下角查找页码文本框：文本为纯数字，位置在页面右下角区域。"""
            for s in get_text_shapes(slide):
                text = s.text_frame.text.strip()
                if not text.isdigit():
                    continue
                x_cm = emu_to_cm(s.left)
                y_cm = emu_to_cm(s.top)
                if x_cm > SLIDE_WIDTH_CM * 0.8 and y_cm > SLIDE_HEIGHT_CM * 0.8:
                    return s
            return None

        def page_number_covers_content(pn_shape, slide):
            """检查页码文本框是否遮挡页面其他内容。
            跳过以下背景底板：
              1. 全页背景（面积 ≥ 80% 页面）
              2. 横贯全宽的页脚底板（宽 ≥ 页面宽 80%，高 ≤ 页面高 20%，且顶部在页面下半）
            其余shape若与页码重叠且页码在其上层则判定为遮挡。"""
            shapes_list = list(slide.shapes)
            pn_z = shapes_list.index(pn_shape) if pn_shape in shapes_list else -1
            slide_area = slide_w * slide_h
            for s in slide.shapes:
                if s is pn_shape:
                    continue
                # 跳过全页背景底板（≥80% 页面面积）
                s_area = s.width * s.height
                if s_area >= slide_area * 0.8:
                    continue
                # 跳过横贯全宽的页脚装饰底板
                if (s.width >= slide_w * 0.8 and
                        s.height <= slide_h * 0.2 and
                        s.top >= slide_h * 0.7):
                    continue
                if not shape_overlaps(pn_shape, s):
                    continue
                s_z = shapes_list.index(s) if s in shapes_list else -1
                if pn_z > s_z:
                    return True
            return False

        failed_pages = []
        covered_pages = []
        for pg in need_pages:
            slide = slides[pg - 1]
            pn = find_page_number_shape(slide)
            if pn is None:
                failed_pages.append(pg)
            elif page_number_covers_content(pn, slide):
                covered_pages.append(pg)

        if failed_pages:
            return False, f"以下页面右下角缺少页码文本: {failed_pages}"
        if covered_pages:
            return False, f"以下页面页码遮挡了其他内容: {covered_pages}"
        return True, ""

    passed_pn, msg_pn = check_page_numbers()
    results.append((5, "需要页码的页面(2,4,5,7,8,10,12,13,15,16)右下角有页码文本且不遮挡内容", passed_pn, msg_pn))

    # --- +3: 第16页新增对象可编辑状态 ---
    def check_slide16_editable():
        """第16页新增对象可编辑状态：
        1. 视频卡片（视频对应的卡片/装饰对象）可单独选中编辑——卡片以独立shape形式存在于页面中。
        2. 页码为文本对象（TEXT_BOX 或 AUTO_SHAPE），不是图片。
        3. 渐变背景可单独选中编辑——以独立shape或页面背景存在，不是和视频合并的整页图片。
        4. 三个视频对象均可单独选中编辑——每个视频为独立的 MEDIA 类型 shape，
           不是和背景合并成整页图片。"""
        slide16 = slides[15]

        # --- 条件3：渐变背景可单独选中编辑 ---
        # 渐变背景需要以"可编辑对象"形式存在，且未与视频合并为整页图片。
        # OOXML 中两种合法写法：
        #   (a) 页面级背景 <p:cSld><p:bg><p:bgPr><a:gradFill/></p:bgPr></p:bg>
        #       在 PowerPoint / WPS 里以"背景"呈现，可通过右键"设置背景格式"编辑；
        #   (b) 独立形状（<p:sp>）覆盖整页且填充 <a:gradFill>：在办公软件里可
        #       直接单击选中作为一个独立 shape 编辑。
        # 反例（不合法）：
        #   - 整页覆盖的 <p:pic>（图片背景）—— 已在"新增视频页"里禁止；此处
        #     仍作把关，若同页存在整页图片且没有独立 gradFill 载体，则说明
        #     渐变背景被并入图片，无法单独编辑；
        #   - 视频对象位于该整页图片之上覆盖 90% 以上 —— 说明视频与背景被
        #     合并成一张整页图片（视频与图片同层）。
        has_bg_gradfill = False
        bg_pr = slide16._element.find('./p:cSld/p:bg/p:bgPr', NS)
        if bg_pr is not None and bg_pr.find('./a:gradFill', NS) is not None:
            has_bg_gradfill = True

        has_shape_gradfill_cover = False
        gradfill_shape = None
        for s in slide16.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.MEDIA:
                continue
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
            if not shape_covers_slide(s):
                continue
            grad_fill = s._element.find('./p:spPr/a:gradFill', NS)
            if grad_fill is not None:
                has_shape_gradfill_cover = True
                gradfill_shape = s
                break

        if not (has_bg_gradfill or has_shape_gradfill_cover):
            return False, ("第16页渐变背景未以页面背景或独立shape形式存在，"
                           "疑似与视频合并成整页图片")

        # 若渐变以独立 shape 存在，进一步校验：视频不应落在该 shape 之上
        # 且该 shape 不应是被锁定/组合进图片的一部分。
        # python-pptx 只暴露 shape 本身，未与背景组合成图片即可视为可独立选中。
        if has_shape_gradfill_cover and gradfill_shape is not None:
            # 若渐变 shape 的父节点是 <p:grpSp>（成组），检查是否与视频同组：
            # 同组则会以整组形式移动，仍视为可单独进入组内编辑（PowerPoint 允许
            # 进入组内选中子对象），此处不判失败；仅当 gradfill_shape 是 PICTURE
            # 才会失败——这不会发生，因为上面已在筛选时排除 PICTURE。
            pass

        # 附加把关：页面中若存在整页 PICTURE，且没有 bg 级 gradFill 承载，
        # 说明"渐变背景"实际由该整页图片提供 —— 视为与背景合并的不合规写法。
        has_full_page_picture = any(
            s.shape_type == MSO_SHAPE_TYPE.PICTURE and shape_covers_slide(s)
            for s in slide16.shapes
        )
        if has_full_page_picture and not has_bg_gradfill:
            return False, "第16页存在整页图片背景，视频与渐变背景疑似合并为整页图片"

        # --- 条件4：三个视频为独立 MEDIA 对象，不是图片 ---
        media = get_media_shapes(slide16)
        if len(media) < 3:
            return False, f"第16页仅有{len(media)}个媒体对象，三个视频对象不足"
        for m in media:
            if m.shape_type != MSO_SHAPE_TYPE.MEDIA:
                return False, f"'{m.name}'不是媒体类型对象，可能已变为图片"

        # 视频未和背景合并为整页图片：若三个视频仍是独立 MEDIA shape，则办公
        # 软件中可单独选中编辑；再叠加上面对整页图片的排查，双重防合并。

        # --- 条件2：页码为文本对象，不是图片 ---
        SLIDE_WIDTH_CM_val = emu_to_cm(slide_w)
        SLIDE_HEIGHT_CM_val = emu_to_cm(slide_h)
        pn_shape = None
        for s in get_text_shapes(slide16):
            text = s.text_frame.text.strip()
            if text.isdigit():
                x_cm = emu_to_cm(s.left)
                y_cm = emu_to_cm(s.top)
                if x_cm > SLIDE_WIDTH_CM_val * 0.8 and y_cm > SLIDE_HEIGHT_CM_val * 0.8:
                    pn_shape = s
                    break
        if pn_shape is None:
            return False, "页码不是文本对象（未在右下角找到纯数字文本框）"
        if pn_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return False, "页码对象是图片，不是文本对象"

        # --- 条件1：视频卡片作为独立shape存在，可单独选中 ---
        # 卡片定义：与视频对象空间重叠或同区域的非视频、非全页背景独立shape
        def related_card_exists(video_shape):
            for s in slide16.shapes:
                if s is video_shape:
                    continue
                if s.shape_type == MSO_SHAPE_TYPE.MEDIA:
                    continue
                if s.width >= slide_w * 0.9 and s.height >= slide_h * 0.9:
                    continue
                # 判断shape中心是否在视频区域内或视频中心是否在shape区域内
                vc_x = video_shape.left + video_shape.width / 2
                vc_y = video_shape.top + video_shape.height / 2
                sc_x = s.left + s.width / 2
                sc_y = s.top + s.height / 2
                in_vid = (video_shape.left <= sc_x <= video_shape.left + video_shape.width and
                          video_shape.top <= sc_y <= video_shape.top + video_shape.height)
                in_card = (s.left <= vc_x <= s.left + s.width and
                           s.top <= vc_y <= s.top + s.height)
                if in_vid or in_card:
                    return True
            return False

        media_sorted = sorted(media, key=lambda s: s.left)
        for m in media_sorted:
            if not related_card_exists(m):
                return False, f"视频'{m.name}'区域未找到对应的独立卡片shape"

        return True, ""

    passed_edit, msg_edit = check_slide16_editable()
    results.append((3, "第16页新增对象可编辑(视频+卡片独立+页码文本+背景独立)", passed_edit, msg_edit))

    # ========== 扣分项 ==========

    # -3: 第1页到第17页中任意一页出现文字超出页面边界、视频遮挡标题、图片遮挡正文、页码遮挡正文、结束页不在最后一页

    def check_spec_layout_issues():
        """严格按细则检查以下5点，不额外约束细则未要求的排版问题：
        1. 第1-17页任意文字超出页面边界；
        2. 第1-17页任意视频对象遮挡标题；
        3. 第1-17页任意图片对象遮挡正文；
        4. 第1-17页任意页码遮挡正文；
        5. 结束页不在最后一页。
        """

        def shape_text(shape):
            return shape.text_frame.text.strip() if shape.has_text_frame else ''

        def overlap_area(s1, s2):
            left = max(s1.left, s2.left)
            top = max(s1.top, s2.top)
            right = min(s1.left + s1.width, s2.left + s2.width)
            bottom = min(s1.top + s1.height, s2.top + s2.height)
            return max(0, right - left) * max(0, bottom - top)

        def has_overlap(s1, s2):
            # Office对象坐标为EMU，给1pt级别容差，避免导出/舍入造成误判。
            return overlap_area(s1, s2) > Pt(1) * Pt(1)

        def is_out_of_slide(shape):
            # 只允许极小Office坐标舍入误差；不能用厘米级容忍，否则会漏判真正越界。
            tolerance = Pt(1)
            return (shape.left < -tolerance or
                    shape.top < -tolerance or
                    shape.left + shape.width > slide_w + tolerance or
                    shape.top + shape.height > slide_h + tolerance)

        def is_footer_shape(shape):
            return bool(shape_text(shape)) and shape.top >= slide_h * 0.86

        def is_page_number(shape, page_no):
            text = shape_text(shape)
            return text in {str(page_no), f'{page_no:02d}'}

        def is_title_shape(shape, page_no):
            text = shape_text(shape)
            if not text or is_footer_shape(shape) or is_page_number(shape, page_no):
                return False
            return shape.top <= slide_h * 0.25 and shape.height >= Cm(0.4)

        def is_body_shape(shape, page_no):
            text = shape_text(shape)
            return (bool(text) and
                    not is_title_shape(shape, page_no) and
                    not is_footer_shape(shape) and
                    not is_page_number(shape, page_no))

        issues = []
        checked_slides = slides[:17]

        for page_no, slide in enumerate(checked_slides, 1):
            shapes = list(slide.shapes)
            text_shapes = [s for s in get_text_shapes(slide) if shape_text(s)]
            picture_shapes = get_picture_shapes(slide)
            media_shapes = get_media_shapes(slide)
            title_shapes = [s for s in text_shapes if is_title_shape(s, page_no)]
            body_shapes = [s for s in text_shapes if is_body_shape(s, page_no)]
            page_number_shapes = [s for s in text_shapes if is_page_number(s, page_no)]

            for text_shape in text_shapes:
                if is_out_of_slide(text_shape):
                    issues.append(f"第{page_no}页文字超出页面边界")
                    break

            for video in media_shapes:
                video_z = shapes.index(video)
                for title in title_shapes:
                    if video_z > shapes.index(title) and has_overlap(video, title):
                        issues.append(f"第{page_no}页视频遮挡标题")
                        break
                else:
                    continue
                break

            for picture in picture_shapes:
                picture_z = shapes.index(picture)
                for body in body_shapes:
                    if picture_z > shapes.index(body) and has_overlap(picture, body):
                        issues.append(f"第{page_no}页图片遮挡正文")
                        break
                else:
                    continue
                break

            for page_number in page_number_shapes:
                for body in body_shapes:
                    if has_overlap(page_number, body):
                        issues.append(f"第{page_no}页页码遮挡正文")
                        break
                else:
                    continue
                break

        end_slide_numbers = [i for i, slide in enumerate(slides, 1) if '感谢观看' in get_all_text(slide)]
        if not end_slide_numbers or end_slide_numbers[-1] != len(slides) or any(i != len(slides) for i in end_slide_numbers):
            issues.append("结束页不在最后一页")

        return len(issues) > 0, '; '.join(issues[:3])

    spec_layout_issue, msg_layout = check_spec_layout_issues()
    results.append((-3, "第1-17页任意页文字超出页面边界/视频遮挡标题/图片遮挡正文/页码遮挡正文/结束页不在最后一页", spec_layout_issue, msg_layout))

    # -3: 第16页视频宽度/高度差异超0.5cm
    def check_video_size_consistency():
        """细则：左侧、中间、右侧三个视频对象的宽度差异超过0.5cm或高度差异超过0.5cm。
        宽度差异 = 三个视频宽度的最大值 - 最小值；高度差异同理。
        任意一项超过0.5cm即触发扣分。"""
        slide16 = slides[15]
        media = sorted(get_media_shapes(slide16), key=lambda s: s.left)
        if len(media) < 3:
            return False, "第16页视频不足3个，无法检查尺寸一致性"
        widths  = [emu_to_cm(m.width)  for m in media]
        heights = [emu_to_cm(m.height) for m in media]
        w_diff = max(widths)  - min(widths)
        h_diff = max(heights) - min(heights)
        if w_diff > 0.5:
            return False, f"三个视频宽度差异{w_diff:.2f}cm，超过0.5cm"
        if h_diff > 0.5:
            return False, f"三个视频高度差异{h_diff:.2f}cm，超过0.5cm"
        return True, ""

    passed_size, msg_size = check_video_size_consistency()
    results.append((-3, "第16页视频宽度/高度差异超过0.5cm", not passed_size, msg_size))

    # -3: 第16页视频重叠或超出页面
    def check_video_overlap_bounds():
        """细则：第16页左侧视频对象、中间视频对象、右侧视频对象之间出现重叠，
        或任意一个视频对象超出页面边界。"""
        slide16 = slides[15]
        media = sorted(get_media_shapes(slide16), key=lambda s: s.left)
        if len(media) < 3:
            return False, "第16页视频不足3个，无法检查左/中/右三个视频对象"

        labels = ["左侧视频对象", "中间视频对象", "右侧视频对象"]
        videos = media[:3]

        # 1. 左侧、中间、右侧三个视频对象之间不能重叠
        for i in range(3):
            for j in range(i + 1, 3):
                if shape_overlaps(videos[i], videos[j]):
                    return False, f"第16页{labels[i]}与{labels[j]}重叠"

        # 2. 任意一个视频对象不能超出页面边界（四边分别检查，基于Office对象边界）
        margin = Cm(0.2)
        for label, video in zip(labels, videos):
            if video.left < -margin:
                return False, f"第16页{label}超出页面左边界"
            if video.top < -margin:
                return False, f"第16页{label}超出页面上边界"
            if video.left + video.width > slide_w + margin:
                return False, f"第16页{label}超出页面右边界"
            if video.top + video.height > slide_h + margin:
                return False, f"第16页{label}超出页面下边界"

        return True, ""

    passed_vob, msg_vob = check_video_overlap_bounds()
    results.append((-3, "第16页左/中/右视频对象重叠或超出页面边界", not passed_vob, msg_vob))

    # -1: 第17页没有出现"感谢观看"文本
    s17_text = get_all_text(slides[16])
    no_thanks = '感谢观看' not in s17_text
    results.append((-1, "第17页没有出现'感谢观看'文本", no_thanks, "第17页未找到'感谢观看'" if no_thanks else ""))

    # -1: 第1、3、6、9、11、14、17页任意一页底部或右下角出现新增页码文本
    def check_no_extra_page_numbers():
        """细则：第1、3、6、9、11、14、17页任意一页底部或右下角出现新增页码文本。
        Office 检查方式：查找这些页面中位于底部区域的纯数字文本对象；
        右下角属于底部区域的子集，因此统一按底部区域检查。"""
        no_num_pages = [1, 3, 6, 9, 11, 14, 17]
        found_pages = []
        for pg in no_num_pages:
            slide = slides[pg - 1]
            for s in get_text_shapes(slide):
                text = s.text_frame.text.strip()
                if not text.isdigit() or int(text) <= 0:
                    continue
                x_cm = emu_to_cm(s.left)
                y_cm = emu_to_cm(s.top)
                # 底部：文本框顶部位于页面高度80%以下；右下角也满足该条件
                in_bottom = y_cm > SLIDE_HEIGHT_CM * 0.8
                in_bottom_right = in_bottom and x_cm > SLIDE_WIDTH_CM * 0.8
                if in_bottom or in_bottom_right:
                    found_pages.append(pg)
                    break
        return found_pages

    extra_pn_pages = check_no_extra_page_numbers()
    results.append((-1, "第1/3/6/9/11/14/17页底部或右下角出现新增页码文本", bool(extra_pn_pages),
                    f"出现页码的页面: {extra_pn_pages}" if extra_pn_pages else ""))

    # -1: 第1页没有出现"HAILAN CARE"文本或没有出现"声波护理品牌汇报"文本
    def check_slide1_required_texts():
        """细则：第1页没有出现'HAILAN CARE'文本或没有出现'声波护理品牌汇报'文本。"""
        s1_text = get_all_text(slides[0])
        missing = []
        if 'HAILAN CARE' not in s1_text:
            missing.append('HAILAN CARE')
        if '声波护理品牌汇报' not in s1_text:
            missing.append('声波护理品牌汇报')
        return len(missing) == 0, f"缺少文本: {missing}" if missing else ""

    s1_text_ok, msg_s1 = check_slide1_required_texts()
    results.append((-1, "第1页没有出现'HAILAN CARE'或'声波护理品牌汇报'文本", not s1_text_ok, msg_s1))

    # -5: 第2页没有出现"目录"文本，或没有出现六组完整目录文本
    def check_page2_toc():
        """细则：第2页必须出现'目录'文本，并出现以下六组完整目录文本：
        '01 品牌概览'、'02 发展路径'、'03 产品生命周期'、
        '04 成熟期策略'、'05 产品组合'、'06 展望未来'。"""
        s2_text = get_all_text(slides[1])
        normalized = re.sub(r'\s+', ' ', s2_text).strip()

        missing = []
        if '目录' not in normalized:
            missing.append('目录')

        toc_items = [
            '01 品牌概览',
            '02 发展路径',
            '03 产品生命周期',
            '04 成熟期策略',
            '05 产品组合',
            '06 展望未来',
        ]
        for item in toc_items:
            if item not in normalized:
                missing.append(item)

        if missing:
            return True, f"缺少文本: {missing}"
        return False, ""

    no_toc, msg_toc = check_page2_toc()
    results.append((-5, "第2页没有出现'目录'或六组完整目录文本", no_toc, msg_toc))

    # -5: 第1-17页中任意一页出现宽度>90%且高度>90%页面的图片，且该页主要中文标题不可单独选中编辑
    def check_fullpage_image_no_editable_title():
        """细则：任意一页出现宽度超过页面宽度90%、高度超过页面高度90%的图片对象，
        且该页主要中文标题不可单独选中编辑（即页面内无独立可编辑的中文标题文本对象）。
        Office 中可单独选中编辑 = 存在独立的文本框 shape，且包含主要中文标题内容。"""
        for i, slide in enumerate(slides, 1):
            pics = get_picture_shapes(slide)
            has_fullpage_pic = any(
                p.width > slide_w * 0.9 and p.height > slide_h * 0.9
                for p in pics
            )
            if not has_fullpage_pic:
                continue

            # 检查该页是否存在独立的、包含中文内容的主要标题文本对象
            has_cn_title = False
            for s in get_text_shapes(slide):
                text = s.text_frame.text.strip()
                # 主要中文标题：含中文字符、长度≥2，且不是页脚（顶部不低于页面85%以下）
                if (len(text) >= 2
                        and any('一' <= c <= '鿿' for c in text)
                        and s.top < slide_h * 0.85):
                    has_cn_title = True
                    break

            if not has_cn_title:
                return True, f"第{i}页有覆盖90%以上页面的图片，但无可单独选中编辑的中文标题文本对象"

        return False, ""

    fullpage_issue, msg_fp = check_fullpage_image_no_editable_title()
    results.append((-5, "任意页面有>90%宽高图片且主要中文标题不可单独选中编辑", fullpage_issue, msg_fp))


    # ========== 汇总结果 ==========
    dim2_items = []
    total_score = 0
    max_score = 0
    for points, desc, triggered, msg in results:
        hit = bool(triggered)
        delta = points if hit else 0
        dim2_items.append({
            "rule": desc,
            "max_delta": points,
            "delta": delta,
            "hit": hit,
            "detail": "",
        })
        total_score += delta
        if points > 0:
            max_score += points

    return {
        "id": script_id,
        "file_name": file_name,
        "status": "ok",
        "error": None,
        "dim1_pass": True,
        "dim1_reason": "",
        "dim2_items": dim2_items,
        "total_score": total_score,
        "max_score": max_score,
    }


if __name__ == "__main__":
    # 本地调试入口：默认以脚本所在目录作为 dir_path
    _here = os.path.dirname(os.path.abspath(__file__))
    _target = sys.argv[1] if len(sys.argv) > 1 else _here
    _report = evaluate(_target)
    print(json.dumps(_report, ensure_ascii=False, indent=2))
