# -*- coding: utf-8 -*-
"""
自动评估脚本：区域协同汇报_第8页Excel超链接版.pptx

评估逻辑：
  - 维度1（可用与可修改性）为门槛维度，不满足直接判 0 分，不再检查维度2。
  - 维度2（完成度）分为加分点与扣分点：
      * 加分细则：必须满足该细则内的每一个点，才累加该条分数。
      * 扣分细则：只要满足该细则内任意一点，即扣分。
  - 最终分数 = 维度2 命中的所有条目分数之和（可正可负）。

依赖：python-pptx
对外接口：evaluate(dir_path: str) -> dict
  - dir_path：脚本所在目录路径；脚本自行在该目录内定位被评估的 .pptx 文件。
  - 返回结构：见评估脚本接口统一约定 §2.2。
"""

import os
import sys
import json
import zipfile
import xml.etree.ElementTree as ET

from pptx import Presentation

SCRIPT_ID = "078"
PPTX_FILE = "区域协同汇报_第8页Excel超链接版.pptx"

# OOXML 命名空间
NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}

EMU_PER_CM = 360000.0


def emu_to_cm(v):
    return float(v) / EMU_PER_CM


# ---------------------------------------------------------------------------
# 几何辅助
# ---------------------------------------------------------------------------
def rects_overlap(a, b):
    """两个矩形 (x, y, w, h) 是否有重叠区域。"""
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return not (ax + aw <= bx or bx + bw <= ax or ay + ah <= by or by + bh <= ay)


def overlap_area(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0, min(ay + ah, by + bh) - max(ay, by))
    return ix * iy


# ---------------------------------------------------------------------------
# 解析第8页：提取形状、超链接、表格、OLE对象等
# ---------------------------------------------------------------------------
class Slide8Data:
    def __init__(self, pptx_path):
        self.pptx_path = pptx_path
        self.slide_count = 0
        self.slide_w_emu = 0
        self.slide_h_emu = 0
        self.shapes = []          # 普通形状 (含超链接信息)
        self.tables = []          # 表格 graphicFrame 区域
        self.ole_objects = []     # 内嵌 OLE 对象区域
        self.title_regions = []   # 顶部标题区域
        self.oval_regions = []    # 椭圆（含底部橙色圆形）区域
        self.rels = {}            # rId -> (Target, TargetMode, Type)
        self._parse()

    def _parse(self):
        prs = Presentation(self.pptx_path)
        self.slide_count = len(prs.slides._sldIdLst)
        self.slide_w_emu = prs.slide_width
        self.slide_h_emu = prs.slide_height

        with zipfile.ZipFile(self.pptx_path) as z:
            slide_xml = z.read("ppt/slides/slide8.xml")
            rels_xml = z.read("ppt/slides/_rels/slide8.xml.rels")

        # 解析 rels
        rels_root = ET.fromstring(rels_xml)
        for rel in rels_root.findall("rel:Relationship", NS):
            rid = rel.get("Id")
            self.rels[rid] = {
                "Target": rel.get("Target"),
                "TargetMode": rel.get("TargetMode", "Internal"),
                "Type": rel.get("Type", ""),
            }

        root = ET.fromstring(slide_xml)
        sptree = root.find(".//p:cSld/p:spTree", NS)

        for el in sptree:
            tag = el.tag.split("}")[-1]
            if tag == "sp":
                self._parse_sp(el)
            elif tag == "pic":
                self._parse_generic_geom(el, "pic")
            elif tag == "graphicFrame":
                self._parse_graphic_frame(el)

    def _get_xfrm(self, spPr_or_frame, frame=False):
        if frame:
            xfrm = spPr_or_frame.find("p:xfrm", NS)
        else:
            xfrm = spPr_or_frame.find("a:xfrm", NS)
        if xfrm is None:
            return None
        off = xfrm.find("a:off", NS)
        ext = xfrm.find("a:ext", NS)
        if off is None or ext is None:
            return None
        return (
            int(off.get("x")),
            int(off.get("y")),
            int(ext.get("cx")),
            int(ext.get("cy")),
        )

    def _gather_text(self, el):
        return "".join(t.text or "" for t in el.iter("{%s}t" % NS["a"]))

    def _parse_sp(self, sp):
        nv = sp.find("p:nvSpPr/p:cNvPr", NS)
        name = nv.get("name") if nv is not None else ""
        spPr = sp.find("p:spPr", NS)
        rect = self._get_xfrm(spPr) if spPr is not None else None

        # 几何类型
        prst = None
        if spPr is not None:
            geom = spPr.find("a:prstGeom", NS)
            if geom is not None:
                prst = geom.get("prst")

        # 超链接：cNvPr 下的 hlinkClick / hlinkHover，或 run 级别的 hlinkClick
        cnvpr_hlink = nv.find("a:hlinkClick", NS) if nv is not None else None
        run_hlinks = sp.findall(".//a:r/a:rPr/a:hlinkClick", NS)

        text = self._gather_text(sp)

        shape_info = {
            "name": name,
            "rect": rect,
            "prst": prst,
            "text": text,
            "shape_level_hlink": cnvpr_hlink.get("{%s}id" % NS["r"]) if cnvpr_hlink is not None else None,
            "run_hlink_ids": [h.get("{%s}id" % NS["r"]) for h in run_hlinks],
            "has_run_hlink_only": (cnvpr_hlink is None and len(run_hlinks) > 0),
        }
        self.shapes.append(shape_info)

        # 椭圆收集（用于底部橙色圆形 / 顶部装饰圆）
        if prst == "ellipse" and rect is not None:
            self.oval_regions.append({"name": name, "rect": rect})

    def _parse_generic_geom(self, el, kind):
        spPr = el.find("p:spPr", NS)
        rect = self._get_xfrm(spPr) if spPr is not None else None
        nv = el.find("p:nvPicPr/p:cNvPr", NS)
        name = nv.get("name") if nv is not None else ""
        descr = nv.get("descr") if nv is not None else ""

        # 图片自身的超链接（cNvPr 下 hlinkClick）
        cnvpr_hlink = nv.find("a:hlinkClick", NS) if nv is not None else None
        shape_level_hlink = (
            cnvpr_hlink.get("{%s}id" % NS["r"]) if cnvpr_hlink is not None else None
        )

        # 图片引用的媒体文件名（blipFill/blip 的 r:embed 对应的 rels Target），
        # 用于识别 Excel/文件图标的语义。
        embed_target = ""
        blip = el.find(".//p:blipFill/a:blip", NS)
        if blip is not None:
            embed_rid = blip.get("{%s}embed" % NS["r"])
            if embed_rid and embed_rid in self.rels:
                embed_target = self.rels[embed_rid].get("Target", "") or ""

        self.shapes.append({
            "name": name, "rect": rect, "prst": None, "text": "",
            "shape_level_hlink": shape_level_hlink,
            "run_hlink_ids": [],
            "has_run_hlink_only": False,
            "kind": kind,
            "descr": descr,
            "embed_target": embed_target,
        })

    def _parse_graphic_frame(self, gf):
        rect = self._get_xfrm(gf, frame=True)
        gdata = gf.find(".//a:graphic/a:graphicData", NS)
        uri = gdata.get("uri") if gdata is not None else ""
        if "table" in uri:
            self.tables.append({"rect": rect})
        elif "ole" in uri:
            # OLE 嵌入对象（Excel）
            ole = gdata.find("p:oleObj", NS)
            rid = ole.get("{%s}id" % NS["r"]) if ole is not None else None
            self.ole_objects.append({"rect": rect, "rid": rid})


# ---------------------------------------------------------------------------
# 评估器
# ---------------------------------------------------------------------------
class Evaluator:
    def __init__(self, data: Slide8Data):
        self.d = data
        self.hits = []          # 命中的细则描述列表
        self.score = 0
        self.page_w_cm = emu_to_cm(data.slide_w_emu)
        self.page_h_cm = emu_to_cm(data.slide_h_emu)
        # 识别主要的"超链接对象"（页面右侧的可点击对象）
        self.link_shape = self._find_link_shape()

    # ---------- 工具 ----------
    def _find_link_shape(self):
        """找出第8页指向外部文件的、位于右侧的超链接形状对象。"""
        candidates = []
        for s in self.d.shapes:
            if s.get("rect") is None:
                continue
            rid = s["shape_level_hlink"] or (s["run_hlink_ids"][0] if s["run_hlink_ids"] else None)
            if rid is None:
                continue
            rel = self.d.rels.get(rid, {})
            candidates.append((s, rid, rel))
        if not candidates:
            return None
        # 优先选择 target 指向 xlsx 的；否则取第一个
        for s, rid, rel in candidates:
            tgt = (rel.get("Target") or "").lower()
            if tgt.endswith(".xlsx") or tgt.endswith(".xls") or "excel" in tgt:
                return {"shape": s, "rid": rid, "rel": rel}
        s, rid, rel = candidates[0]
        return {"shape": s, "rid": rid, "rel": rel}

    def _link_rect_cm(self):
        if not self.link_shape:
            return None
        x, y, w, h = self.link_shape["shape"]["rect"]
        return (emu_to_cm(x), emu_to_cm(y), emu_to_cm(w), emu_to_cm(h))

    # =====================================================================
    # 维度1：可用与可修改性（门槛）
    # =====================================================================
    def check_dimension1(self):
        reasons = []
        ok = True

        # 1) 格式为 .pptx（python-pptx 无法可靠打开二进制 .ppt，此处严格要求 .pptx）
        ext = os.path.splitext(self.d.pptx_path)[1].lower()
        if ext != ".pptx":
            ok = False
            reasons.append(f"文件格式 {ext} 不是 .pptx")
            return ok, reasons
        reasons.append(f"格式 {ext} 合法")

        # 2) 能正常打开：.pptx 必须是有效 OOXML(ZIP) 包，且核心部件齐全
        can_open = False
        try:
            with zipfile.ZipFile(self.d.pptx_path) as z:
                names = set(z.namelist())
                required = {"[Content_Types].xml", "ppt/presentation.xml"}
                missing = required - names
                if missing:
                    ok = False
                    reasons.append(f"缺少 .pptx 核心部件：{sorted(missing)}，无法正常打开")
                else:
                    can_open = True
                    reasons.append("文件为有效的 .pptx(OOXML) 包，核心部件齐全，可正常打开")
        except zipfile.BadZipFile:
            ok = False
            reasons.append("文件不是有效的 .pptx(ZIP) 包，无法正常打开")
        except Exception as exc:
            ok = False
            reasons.append(f"打开文件失败：{type(exc).__name__}: {exc}")

        if not can_open:
            return ok, reasons

        # 3) 保持 10 页
        if self.d.slide_count != 10:
            ok = False
            reasons.append(f"页数为 {self.d.slide_count}，不是 10 页")
        else:
            reasons.append("PPT 共 10 页")

        return ok, reasons

    # =====================================================================
    # 维度2：完成度
    # =====================================================================
    def _add(self, score, desc):
        self.score += score
        self.hits.append((score, desc))

    # ---- +5：第8页右侧空白区域超链接（位置/尺寸/图标按钮，所有点都需满足）----
    # 细则原文逐点拆解：
    #   1) 左边缘距页面左边线约 26.0 至 29.0 厘米
    #   2) 整体位于页面右侧约 77% 至 95% 区域
    #   3) 上边缘距页面上边线约 5.0 至 11.0 厘米
    #   4) 不进入顶部标题区或底部橙色圆形区域
    #   5) 宽度约 3.5 至 5.5 厘米
    #   6) 高度约 1.5 至 3.5 厘米
    #   7) 能够清楚点击且不过度占据空白区域
    #   8) 使用 Excel 图标、文件图标或带文字的按钮超链接，能够直观表达“打开 Excel 文件”
    def rule_plus5_position(self):
        sub = []          # (满足?, 说明)
        rect = self._link_rect_cm()
        if rect is None:
            return False, ["未找到第8页的超链接对象，无法判定位置/尺寸"]

        x, y, w, h = rect

        # 1) 左边缘距页面左边线约 26.0 至 29.0 厘米
        p1 = 26.0 <= x <= 29.0
        sub.append((p1, f"左边缘距页面左边线 {x:.1f}cm（要求约 26.0~29.0cm）"))

        # 2) 整体位于页面右侧约 77% 至 95% 区域
        #    “整体”按对象左右边缘相对页宽的占位区间判断，须落在 77%~95% 内。
        left_ratio = x / self.page_w_cm
        right_ratio = (x + w) / self.page_w_cm
        p2 = left_ratio >= 0.77 and right_ratio <= 0.95
        sub.append((p2, f"整体位于页面 {left_ratio*100:.0f}%~{right_ratio*100:.0f}% 区域（要求约 77%~95%）"))

        # 3) 上边缘距页面上边线约 5.0 至 11.0 厘米
        p3 = 5.0 <= y <= 11.0
        sub.append((p3, f"上边缘距页面上边线 {y:.1f}cm（要求约 5.0~11.0cm）"))

        # 4) 不进入顶部标题区或底部橙色圆形区域
        link_rect_emu = self.link_shape["shape"]["rect"]
        # 顶部标题区：用页面顶部约 0~3.5cm 的横向条带近似
        title_band = (0, 0, self.d.slide_w_emu, int(3.5 * EMU_PER_CM))
        in_title = rects_overlap(link_rect_emu, title_band)
        # 底部橙色圆形：取位于页面底部区域的椭圆作几何重叠检测
        bottom_ovals = [o for o in self.d.oval_regions
                        if o["rect"] and o["rect"][1] > self.d.slide_h_emu * 0.6]
        in_orange = any(rects_overlap(link_rect_emu, o["rect"]) for o in bottom_ovals)
        p4 = (not in_title) and (not in_orange)
        sub.append((p4, f"未进入顶部标题区(进入={in_title})且未进入底部橙色圆形(进入={in_orange})"))

        # 5) 宽度约 3.5 至 5.5 厘米
        p5 = 3.5 <= w <= 5.5
        sub.append((p5, f"宽度 {w:.1f}cm（要求约 3.5~5.5cm）"))

        # 6) 高度约 1.5 至 3.5 厘米
        p6 = 1.5 <= h <= 3.5
        sub.append((p6, f"高度 {h:.1f}cm（要求约 1.5~3.5cm）"))

        # 7) 能够清楚点击且不过度占据空白区域
        #    “能够清楚点击”：对象本身具备一定面积（不是退化为线/点），按可点击下限约 1.0cm×0.5cm 判断。
        #    “不过度占据空白区域”：对象面积不超过其所在右侧空白区（约页宽 23%）的合理上限，这里取整页面积的 10%。
        clickable = (w >= 1.0 and h >= 0.5)
        obj_area = w * h
        page_area = self.page_w_cm * self.page_h_cm
        not_overcrowd = obj_area <= page_area * 0.10
        p7 = clickable and not_overcrowd
        sub.append((p7, f"可清楚点击(面积 {obj_area:.1f}cm²，占页面 {obj_area/page_area*100:.1f}%，未过度占据空白)"))

        # 8) 使用 Excel 图标、文件图标或带文字的按钮超链接，能够直观表达“打开 Excel 文件”
        shp = self.link_shape["shape"]
        text = (shp.get("text") or "")
        text_l = text.lower()
        # 按钮文字关键词收紧到明确表达“打开 Excel 文件”的词，
        # 移除 crm / pipeline / spreadsheet / 泛化的 file / 文件 等不足以直观表达的词。
        express_kw = ["excel", "xlsx", "xls", "xlsm",
                      "打开", "open",
                      "表格", "工作簿", "工作表"]
        # 带文字的按钮超链接：按钮形状（roundRect/rect）且含可直观表达的文字，
        # 且超链接绑定在形状(按钮)主体上（否则按钮本体不可点击）。
        is_button = shp.get("prst") in ("roundRect", "rect")
        has_text = len(text.strip()) > 0
        expresses_text = any(k in text_l for k in express_kw)
        button_has_hlink = shp.get("shape_level_hlink") is not None
        button_ok = is_button and has_text and expresses_text and button_has_hlink

        # 图标超链接：对象为图片/图标（kind=pic），
        # 图片/图标对象自身需带有超链接（cNvPr/hlinkClick），
        # 且从图片的名称/替代文字(descr)/引用的媒体文件名中能识别出 Excel/文件图标语义。
        is_icon_kind = shp.get("kind") == "pic"
        icon_has_hlink = shp.get("shape_level_hlink") is not None
        icon_meta = " ".join([
            str(shp.get("name") or ""),
            str(shp.get("descr") or ""),
            str(shp.get("embed_target") or ""),
        ]).lower()
        icon_semantic_kw = ["excel", "xlsx", "xls", "xlsm",
                            "workbook", "worksheet", "spreadsheet",
                            "表格", "工作簿", "工作表",
                            "file", "document", "文件", "文档"]
        icon_expresses = any(k in icon_meta for k in icon_semantic_kw)
        icon_ok = is_icon_kind and icon_has_hlink and icon_expresses

        p8 = button_ok or icon_ok
        sub.append((p8, (
            f"使用带文字按钮/Excel或文件图标超链接："
            f"按钮命中={button_ok}(prst={shp.get('prst')},文字=“{text}”,形状级链接={button_has_hlink})，"
            f"图标命中={icon_ok}(kind={shp.get('kind')},自身链接={icon_has_hlink},"
            f"语义源=“{icon_meta[:40]}”)"
        )))

        all_ok = all(s[0] for s in sub)
        detail = [("  ✓ " if s[0] else "  ✗ ") + s[1] for s in sub]
        return all_ok, detail

    # ---- +5：第8页放映状态（点击能调用Excel打开工作簿）----
    # 细则原文逐点拆解：
    #   1) 单击超链接对象后能够调用 Excel 或系统默认表格程序打开工作簿
    #   2) 通过 PowerPoint 支持的“Ctrl+单击”或“打开超链接”操作能够访问 Excel 文件
    def _resolve_workbook_target(self, tgt, mode):
        """解析超链接 Target，返回 (is_accessible, kind, resolved)。

        kind ∈ {"absolute","relative","embedded","url","empty","other"}。
          - absolute / relative：磁盘路径，用 os.path.isfile 判断实际存在。
          - embedded：pptx 包内部件路径（TargetMode=Internal），用 zipfile 校验是否在包内。
          - url：http/https/ftp/mailto/file 协议，不做联网访问，视为不可确定存在——按不可访问处理。
          - empty / other：目标为空或非文件形态，按不可访问处理。
        """
        if not tgt:
            return False, "empty", ""

        tl = tgt.lower()
        # URL / 协议链接：本地无法验证工作簿是否可访问，且 rubric 关注的是打开工作簿而非网页
        if tl.startswith(("http://", "https://", "ftp://", "mailto:")):
            return False, "url", tgt
        if tl.startswith("file:"):
            # 转成本地路径再判文件存在
            local = tgt[5:].lstrip("/")
            local = local.replace("/", os.sep)
            return os.path.isfile(local), "url", local

        pptx_dir = os.path.dirname(os.path.abspath(self.d.pptx_path))
        # 绝对路径：盘符 C:\ / UNC \\server\share / POSIX /path
        norm = tgt.replace("/", os.sep)
        is_abs = (
            (len(tgt) >= 2 and tgt[1] == ":") or
            tgt.startswith("\\\\") or
            norm.startswith(os.sep)
        )
        if mode == "External":
            if is_abs:
                return os.path.isfile(norm), "absolute", norm
            # 相对路径：相对于 .pptx 所在目录解析
            resolved = os.path.normpath(os.path.join(pptx_dir, norm))
            return os.path.isfile(resolved), "relative", resolved

        # TargetMode 为 Internal：视为 pptx 包内嵌入部件（如 ../embeddings/xxx.xlsx）
        # rels Target 是相对 ppt/slides/_rels/slide8.xml.rels 的路径
        try:
            part_path = os.path.normpath(os.path.join("ppt/slides", tgt)).replace("\\", "/")
        except Exception:
            part_path = tgt.replace("\\", "/")
        try:
            with zipfile.ZipFile(self.d.pptx_path) as z:
                exists = part_path in set(z.namelist())
        except Exception:
            exists = False
        return exists, "embedded", part_path

    def rule_plus5_show(self):
        sub = []
        if not self.link_shape:
            return False, ["未找到超链接对象，无法判定放映点击行为"]
        rel = self.link_shape["rel"]
        tgt = (rel.get("Target") or "")
        mode = rel.get("TargetMode", "Internal")
        rtype = rel.get("Type", "")
        tl = tgt.lower()
        shp = self.link_shape["shape"]

        # 目标文件存在性/可访问性解析
        accessible, tgt_kind, tgt_resolved = self._resolve_workbook_target(tgt, mode)
        points_to_workbook_ext = tl.endswith((".xlsx", ".xls", ".xlsm"))

        # 1) 单击超链接对象后能够调用 Excel 或系统默认表格程序打开工作簿
        #    判据：超链接挂在“对象”上（形状级 hlinkClick，整体可单击）；
        #          目标为可由 Excel / 系统默认表格程序打开的工作簿文件（.xlsx/.xls/.xlsm）；
        #          且目标工作簿实际存在/可访问（绝对/相对路径 → 磁盘文件存在；嵌入 → 包内部件存在）。
        on_object = shp.get("shape_level_hlink") is not None
        p1 = on_object and points_to_workbook_ext and accessible
        sub.append((p1, (
            f"单击超链接对象(对象级超链接={on_object})可调用Excel/默认表格程序打开工作簿"
            f"（扩展名匹配={points_to_workbook_ext}，目标类型={tgt_kind}，"
            f"实际可访问={accessible}，解析路径=“{tgt_resolved}”）"
        )))

        # 2) 通过 PowerPoint 支持的“Ctrl+单击”或“打开超链接”操作能够访问 Excel 文件
        #    判据：为 PowerPoint 标准 hyperlink 关系（外部 Target）或嵌入 Excel 的 package 关系；
        #          且目标扩展名为工作簿；且目标实际可访问——只有可访问时才算“能够访问 Excel 文件”。
        is_hyperlink_rel = rtype.endswith("/hyperlink")
        is_package_rel = "/oleObject" in rtype or "/package" in rtype
        is_pp_supported = is_hyperlink_rel or is_package_rel
        p2 = is_pp_supported and points_to_workbook_ext and accessible
        sub.append((p2, (
            f"为PowerPoint支持的超链接/嵌入关系(hyperlink={is_hyperlink_rel},"
            f"package={is_package_rel},外部={mode == 'External'})，"
            f"目标可访问={accessible}，可经“Ctrl+单击/打开超链接”访问Excel文件"
        )))

        all_ok = all(s[0] for s in sub)
        detail = [("  ✓ " if s[0] else "  ✗ ") + s[1] for s in sub]
        return all_ok, detail

    # =====================================================================
    # 主流程
    # =====================================================================
    def run(self):
        """执行评估，返回结构化 dict（详见接口约定 §2.2）。"""
        file_name = os.path.basename(self.d.pptx_path)

        # ---- 维度2 规则清单（预先定义，便于在 dim1 未通过时也能列出 max_score）----
        plus_rules = [
            (5,
             "第8页右侧空白区域超链接：左边缘距页面左边线约26.0至29.0厘米，整体位于页面右侧约77%至95%区域。"
             "上边缘距页面上边线约5.0至11.0厘米，不进入顶部标题区或底部橙色圆形区域。宽度约3.5至5.5厘米、"
             "高度约1.5至3.5厘米，能够清楚点击且不过度占据空白区域。使用Excel图标、文件图标或带文字的按钮超链接，"
             "能够直观表达“打开Excel文件”。",
             self.rule_plus5_position),
            (5,
             "第8页放映状态：单击超链接对象后能够调用Excel或系统默认表格程序打开工作簿。"
             "通过PowerPoint支持的“Ctrl+单击”或“打开超链接”操作能够访问Excel文件。",
             self.rule_plus5_show),
        ]
        # 维度二满分 = 加分项 max_delta 之和
        max_score = sum(sc for sc, _, _ in plus_rules)

        # ---- 维度1（门槛）----
        d1_ok, d1_reasons = self.check_dimension1()
        if not d1_ok:
            return {
                "id": SCRIPT_ID,
                "file_name": file_name,
                "status": "ok",
                "error": None,
                "dim1_pass": False,
                "dim1_reason": "；".join(d1_reasons),
                "dim2_items": [],
                "total_score": 0,
                "max_score": max_score,
            }

        # ---- 维度2：完成度 ----
        dim2_items = []

        # 加分项：需满足细则内每一个点才命中
        for sc, rule_text, fn in plus_rules:
            ok = fn()[0]
            delta = sc if ok else 0
            self.score += delta
            dim2_items.append({
                "rule": rule_text,
                "max_delta": sc,
                "delta": delta,
                "hit": bool(ok),
                "detail": "",
            })

        # 扣分项：满足任意一点即命中（当前无扣分项）

        return {
            "id": SCRIPT_ID,
            "file_name": file_name,
            "status": "ok",
            "error": None,
            "dim1_pass": True,
            "dim1_reason": "",
            "dim2_items": dim2_items,
            "total_score": self.score,
            "max_score": max_score,
        }


def evaluate(dir_path: str) -> dict:
    """
    统一入口：接收脚本所在目录的路径，脚本自己在该目录里定位并打开被评估的文档。

    参数:
        dir_path: 脚本所在目录（其中含被评估的 .pptx 文件）。
    返回:
        结构化 dict，详见接口约定 §2.2。
    """
    file_name = PPTX_FILE
    try:
        # 优先使用约定文件名；若不存在，回退为在目录内查找首个 .pptx 文件
        pptx_path = os.path.join(dir_path, PPTX_FILE)
        if not os.path.isfile(pptx_path):
            candidates = [f for f in os.listdir(dir_path)
                          if f.lower().endswith(".pptx")]
            if not candidates:
                return {
                    "id": SCRIPT_ID,
                    "file_name": PPTX_FILE,
                    "status": "error",
                    "error": f"目录 {dir_path} 下未找到 .pptx 文件",
                    "dim1_pass": False,
                    "dim1_reason": "",
                    "dim2_items": [],
                    "total_score": 0,
                    "max_score": 10,
                }
            pptx_path = os.path.join(dir_path, candidates[0])
            file_name = candidates[0]

        data = Slide8Data(pptx_path)
        return Evaluator(data).run()
    except Exception as exc:
        return {
            "id": SCRIPT_ID,
            "file_name": file_name,
            "status": "error",
            "error": f"{type(exc).__name__}: {exc}",
            "dim1_pass": False,
            "dim1_reason": "",
            "dim2_items": [],
            "total_score": 0,
            "max_score": 10,
        }


if __name__ == "__main__":
    # 本地调试：默认使用脚本所在目录；也可通过命令行参数覆盖目录。
    _dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    _result = evaluate(_dir)
    # 用 ensure_ascii=True 避免非 UTF-8 控制台（如 Windows cp1252）编码错误。
    sys.stdout.write(json.dumps(_result, ensure_ascii=True, indent=2) + "\n")
