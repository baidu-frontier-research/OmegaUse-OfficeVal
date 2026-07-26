#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
客户投诉闭环处理流程 PPT 自动评分脚本。

对外接口：
    evaluate(dir_path: str) -> dict
        入参为“脚本所在目录的路径”；脚本自身在该目录里定位并打开被评估的 .pptx 文件，
        返回统一约定的结构化结果字典（详见函数 docstring）。

评估逻辑：
1. 先做“维度1：可用与可修改性”门槛检查；任一门槛不满足，直接 0 分，不进入维度2。
2. 维度1通过后，对“维度2：完成度评分细则”的每个得分点做自动化检测；命中则累加该点分值。
3. 结果全部以 dict 形式 return，不 print 主结果、不改 sys.stdout、不 sys.exit。

依赖：python-pptx、lxml
安装：pip install python-pptx lxml
本地自测：python officeval_056_verifier.py <脚本所在目录>
"""
from __future__ import annotations

import json
import math
import os
import re
import sys
import zipfile
from dataclasses import dataclass
from typing import Callable, Iterable, Optional

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

EMU_PER_CM = 360000
EMU_PER_PT = 12700

TARGET_TITLE = "客户投诉闭环处理流程"
TOP_BLUE_LABELS = ["接收投诉", "信息登记", "问题分级"]
TOP_GREEN_LABELS = ["是否重大问题？", "客服处理", "客户确认", "结案归档"]
BOTTOM_LABELS = ["指派专员", "原因分析", "制定整改方案", "执行整改", "质量复核", "复核通过？"]
RECORD_LABELS = ["相关记录：", "沟通记录", "证据材料", "处理报告"]
ALL_CORE_LABELS = TOP_BLUE_LABELS + TOP_GREEN_LABELS + BOTTOM_LABELS + RECORD_LABELS

# 颜色按“评分意图”归类，而不是要求完全等于某个 RGB。
# 主题色（深蓝/青绿/橙 及其浅色）按色系（含亮度/通道主导）判定；
# 非主题色（灰/黑/白等中性色）仅按“颜色类别”本身判定——只要看起来是灰就算灰、是白就算白，
# 不受具体明度或特定色值的限制。
# GRAY_MAX 用于把“白色”从“灰/黑”里排除掉：三通道均 >= GRAY_MAX 时判为白，其余中性色都算灰/黑。
GRAY_MAX = 230


def cm(value) -> float:
    return float(value) / EMU_PER_CM if value is not None else 0.0


def pt(value) -> Optional[float]:
    return float(value) / EMU_PER_PT if value is not None else None


def norm_text(text: str) -> str:
    text = text or ""
    text = text.replace("　", " ")
    text = re.sub(r"\s+", "", text)
    text = text.replace("?", "？").replace(":", "：")
    return text


def text_of(shape) -> str:
    if getattr(shape, "has_text_frame", False):
        return shape.text_frame.text or ""
    return ""


def shape_xml(shape) -> str:
    try:
        return etree.tostring(shape.element, encoding="unicode")
    except Exception:
        return ""


def auto_shape_name(shape) -> str:
    try:
        return str(shape.auto_shape_type)
    except Exception:
        return ""


def is_rounded_rect(shape) -> bool:
    return "ROUNDED_RECTANGLE" in auto_shape_name(shape) or "roundRect" in shape_xml(shape)


def is_diamond(shape) -> bool:
    return "DIAMOND" in auto_shape_name(shape) or 'prst="diamond"' in shape_xml(shape)


def is_oval(shape) -> bool:
    return "OVAL" in auto_shape_name(shape) or 'prst="ellipse"' in shape_xml(shape)


def is_rectangle(shape) -> bool:
    # 圆角矩形与直角矩形都属于“矩形框”。
    return "RECTANGLE" in auto_shape_name(shape) or 'prst="rect"' in shape_xml(shape) or "roundRect" in shape_xml(shape)


def is_white_fill(rgb: Optional[tuple[int, int, int]]) -> bool:
    # “白底”：办公软件里呈现为白/近白（各通道都很高）。
    if rgb is None:
        return False
    r, g, b = rgb
    return min(r, g, b) >= 230


def is_line(shape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.LINE or shape.element.tag.endswith("cxnSp")


def is_arrow_shape(shape) -> bool:
    """判断形状是否为“一体箭头”——PowerPoint/WPS 内置的整体箭头自选图形。

    对应 DrawingML 中 <a:prstGeom prst="..."/> 为 *Arrow* 家族的形状，例如
    rightArrow、leftArrow、upArrow、downArrow、bentArrow、curvedArrow、
    leftRightArrow、quadArrow 等。这些是“一支完整箭头”的独立可编辑图形，
    与“线条+箭头端点”表现等价，因此在箭头判定中同样应算作合法箭头对象。
    """
    xml = shape_xml(shape)
    return bool(re.search(r'prst="\w*[Aa]rrow\w*"', xml))


def is_arrow_object(shape) -> bool:
    """“箭头对象”统一判定：line / 直线连接符 / 一体箭头 均算箭头对象。"""
    return is_line(shape) or is_arrow_shape(shape)


def arrow_rgb(shape) -> Optional[tuple[int, int, int]]:
    """获取箭头对象的显示颜色。

    - 线条/连接线：颜色由 <a:ln> 描述，取线条颜色；
    - 一体箭头图形：颜色由整体形状的填充呈现，取填充色 fore_color；填充缺失
      时退回描边色，兼容部分模板把颜色只写在描边上的情况。
    """
    if is_arrow_shape(shape) and not is_line(shape):
        c = fill_rgb(shape)
        if c is not None:
            return c
    return line_rgb(shape)


def is_picture(shape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE


def is_group(shape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.GROUP


def iter_shapes(slide_or_group) -> Iterable:
    for shape in slide_or_group.shapes:
        yield shape
        if is_group(shape):
            try:
                yield from iter_shapes(shape)
            except Exception:
                pass


def rgb_from_color(color) -> Optional[tuple[int, int, int]]:
    try:
        rgb = color.rgb
        if rgb is None:
            return None
        s = str(rgb).upper()
        return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
    except Exception:
        return None


def line_rgb(shape) -> Optional[tuple[int, int, int]]:
    try:
        return rgb_from_color(shape.line.color)
    except Exception:
        return None


def fill_rgb(shape) -> Optional[tuple[int, int, int]]:
    try:
        return rgb_from_color(shape.fill.fore_color)
    except Exception:
        return None


def run_rgb(run) -> Optional[tuple[int, int, int]]:
    try:
        return rgb_from_color(run.font.color)
    except Exception:
        return None


def rgb_hex(rgb: Optional[tuple[int, int, int]]) -> str:
    if rgb is None:
        return "None"
    return "%02X%02X%02X" % rgb


def color_distance(a: tuple[int, int, int], b: tuple[int, int, int]) -> float:
    return math.sqrt(sum((x - y) ** 2 for x, y in zip(a, b)))


def is_dark_blue(rgb: Optional[tuple[int, int, int]]) -> bool:
    if rgb is None:
        return False
    r, g, b = rgb
    samples = [(9, 54, 122), (0, 79, 185), (0, 63, 150)]
    return any(color_distance(rgb, s) <= 85 for s in samples) or (b >= 100 and b > r * 1.6 and b > g * 1.15)


def is_blue(rgb: Optional[tuple[int, int, int]]) -> bool:
    if rgb is None:
        return False
    r, g, b = rgb
    return is_dark_blue(rgb) or (b >= 120 and b > r and b >= g * 0.8)


def is_teal_or_green(rgb: Optional[tuple[int, int, int]]) -> bool:
    if rgb is None:
        return False
    r, g, b = rgb
    samples = [(0, 145, 135), (0, 128, 110), (0, 120, 80), (0, 150, 90)]
    return any(color_distance(rgb, s) <= 90 for s in samples) or (g >= 100 and g > r * 1.5 and b >= 60)


def is_orange(rgb: Optional[tuple[int, int, int]]) -> bool:
    if rgb is None:
        return False
    r, g, b = rgb
    samples = [(235, 111, 0), (237, 125, 49), (255, 153, 51)]
    return any(color_distance(rgb, s) <= 85 for s in samples) or (r >= 180 and 60 <= g <= 180 and b <= 90)


def is_light_blue(rgb: Optional[tuple[int, int, int]]) -> bool:
    if rgb is None:
        return False
    r, g, b = rgb
    return r >= 210 and g >= 225 and b >= 235 and b >= r


def is_light_green(rgb: Optional[tuple[int, int, int]]) -> bool:
    if rgb is None:
        return False
    r, g, b = rgb
    return r >= 220 and g >= 235 and b >= 220 and g >= r - 10


def is_light_orange(rgb: Optional[tuple[int, int, int]]) -> bool:
    if rgb is None:
        return False
    r, g, b = rgb
    return r >= 240 and g >= 225 and b >= 220 and r >= g - 20


def is_black_or_dark_gray(rgb: Optional[tuple[int, int, int]]) -> bool:
    if rgb is None:
        # 部分文本继承主题色，无法直接取 RGB 时按不通过处理，避免误判。
        return False
    r, g, b = rgb
    # 非主题色（灰/黑）判定：只要三通道接近（视觉上为中性灰），且不属于白色即可。
    # 不再限制“暗到什么程度”——只要是灰/黑色系（含中灰），就算通过。
    neutral = abs(r - g) <= 35 and abs(g - b) <= 35 and abs(r - b) <= 35
    not_white = min(r, g, b) < GRAY_MAX
    return neutral and not_white


def all_runs(shape):
    if not getattr(shape, "has_text_frame", False):
        return []
    runs = []
    for p in shape.text_frame.paragraphs:
        runs.extend(p.runs)
    return runs


def text_font_sizes(shape) -> list[float]:
    sizes = []
    for r in all_runs(shape):
        if r.text.strip() and r.font.size is not None:
            sizes.append(pt(r.font.size))
    return sizes


def text_colors(shape) -> list[Optional[tuple[int, int, int]]]:
    return [run_rgb(r) for r in all_runs(shape) if r.text.strip()]


def text_font_names(shape) -> list[str]:
    names = []
    for r in all_runs(shape):
        if not r.text.strip():
            continue
        name = r.font.name
        if name:
            names.append(name)
            continue
        xml = shape_xml(shape)
        m = re.search(r'typeface="([^"]+)"', xml)
        names.append(m.group(1) if m else "")
    return names


def run_typefaces(run) -> dict:
    """读取单个 run 自身设置的字体（只看这个 run，不跨 run 取值）。

    办公软件（PowerPoint/WPS）在渲染一个 run 时按字符所属文种分别选字体：
    - 中文/东亚字符 -> ea（<a:ea>）
    - 拉丁字母、数字 -> latin（<a:latin>）
    - 复杂文种 -> cs（<a:cs>）
    因此判断“文字在办公软件里实际显示成什么字体”，必须看对应槽位，而不能只看 latin。
    """
    result = {"latin": None, "ea": None, "cs": None}
    try:
        xml = etree.tostring(run._r, encoding="unicode")
    except Exception:
        return result
    for slot in ("latin", "ea", "cs"):
        m = re.search(r"<a:%s\b[^>]*\btypeface=\"([^\"]*)\"" % slot, xml)
        if m:
            result[slot] = m.group(1)
    return result


def is_source_handwriting_font(name: Optional[str]) -> bool:
    return bool(name) and "思源手写" in name


def has_chinese(text: str) -> bool:
    return re.search(r"[一-鿿]", text or "") is not None


def paragraph_align_center(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    paragraphs = [p for p in shape.text_frame.paragraphs if p.text.strip()]
    return bool(paragraphs) and all(p.alignment in (PP_ALIGN.CENTER, None) for p in paragraphs)


def paragraph_align_left(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    paragraphs = [p for p in shape.text_frame.paragraphs if p.text.strip()]
    return bool(paragraphs) and all(p.alignment in (PP_ALIGN.LEFT, None) for p in paragraphs)


def vertical_middle(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    return shape.text_frame.vertical_anchor in (MSO_ANCHOR.MIDDLE, None)


def vertical_top(shape) -> bool:
    if not getattr(shape, "has_text_frame", False):
        return False
    return shape.text_frame.vertical_anchor in (MSO_ANCHOR.TOP, None)


def line_width_pt(shape) -> Optional[float]:
    try:
        if shape.line.width is None:
            return None
        return pt(shape.line.width)
    except Exception:
        return None


# PowerPoint/WPS 中未显式设置 <a:ln w=".."/> 时，形状轮廓按默认线宽 1.0 磅渲染。
# 因此在线宽判定时，把“未设置”视为默认值 1.0pt，只要该默认值落在要求范围内即通过，
# 与办公软件的真实渲染一致。
DEFAULT_LINE_WIDTH_PT = 1.0


def effective_line_width_pt(shape) -> float:
    """获取形状用于判定的“有效线宽”（磅）。

    - 若显式设置了线宽，则返回该值；
    - 未显式设置时，返回办公软件默认线宽 1.0pt（对应实际渲染出的边线粗细）。
    """
    w = line_width_pt(shape)
    return DEFAULT_LINE_WIDTH_PT if w is None else w


def has_edge_line(shape) -> bool:
    """判断形状轮廓是否为“边线”（描边），对应办公软件里可见的框线。

    DrawingML 中形状描边由 <a:ln> 描述；若其内为 <a:noFill> 则表示“无线条”。
    有些形状不写显式 <a:ln>，而通过样式引用 <a:lnRef> 套用主题描边，此时
    以能取到线条颜色作为“存在边线”的证据。
    """
    xml = shape_xml(shape)
    m = re.search(r"<a:ln\b.*?</a:ln>", xml, re.S)
    if m:
        return "<a:noFill" not in m.group(0)
    if re.search(r"<a:ln\b[^>]*/>", xml):
        return True
    return bool(re.search(r"<a:lnRef\b", xml)) and line_rgb(shape) is not None


def has_arrowhead(shape) -> bool:
    xml = shape_xml(shape)
    # 一体箭头（预设几何为 *Arrow*）本身就是一支完整箭头，视为“带箭头”。
    if is_arrow_shape(shape):
        return True
    return "tailEnd" in xml or "headEnd" in xml


def arrow_end_types(shape) -> tuple[Optional[str], Optional[str]]:
    """读取连接线两端的箭头类型 (headEnd, tailEnd)。

    DrawingML 中 <a:ln> 下的 <a:headEnd type="..."/> / <a:tailEnd type="..."/>
    描述线段起点/终点的端点样式。type 为 "none" 或缺失表示该端无箭头。
    """
    xml = shape_xml(shape)
    h = re.search(r'<a:headEnd\b[^>]*\btype="([^"]+)"', xml)
    t = re.search(r'<a:tailEnd\b[^>]*\btype="([^"]+)"', xml)
    ht = h.group(1) if h else None
    tt = t.group(1) if t else None
    ht = None if ht in (None, "none") else ht
    tt = None if tt in (None, "none") else tt
    return ht, tt


def is_single_lr_arrow(shape) -> bool:
    """判断形状是否为“单向、方向从左向右”的箭头（在办公软件中真实生效的方向）。

    两种常见画法都覆盖：
    1) 预设箭头形状：prst=rightArrow 指向右；leftArrow 需 flipH 才指向右；
       leftRightArrow/其它双向箭头不算单向。
    2) 连接线端点：连接线默认从局部左端(headEnd)指向右端(tailEnd)，flipH 会左右镜像。
       右端有箭头 = 办公软件中箭头指向右：
         (tailEnd 为箭头 且 未翻转) 或 (headEnd 为箭头 且 翻转)。
       单向 = 仅一端有箭头。
    """
    xml = shape_xml(shape)
    flip = 'flipH="1"' in xml
    m = re.search(r'prst="(\w*[Aa]rrow\w*)"', xml)
    if m:
        prst = m.group(1).lower()
        if "leftright" in prst or "quad" in prst:
            return False
        if "right" in prst:
            return not flip
        if "left" in prst:
            return flip
    head, tail = arrow_end_types(shape)
    n_arrows = (1 if head else 0) + (1 if tail else 0)
    if n_arrows != 1:
        return False
    return (tail is not None and not flip) or (head is not None and flip)


def arrow_points_down(shape) -> bool:
    """判断箭头对象的实际指向是否“向下”（在办公软件中呈现的方向）。

    与 is_single_lr_arrow 同思路，覆盖两种画法：
    1) 一体箭头：prst=downArrow 默认向下，flipV 后转为向上；prst=upArrow 反之。
       leftRightArrow/quadArrow/bentArrow/curvedArrow 等非单向或方向不定的箭头，
       无法从预设几何稳定判断“向下”，返回 False。
    2) 连接线：bbox 的对角线两端即两个端点，未 flipV 时从 (x1,y1) 到 (x2,y2)——
       即 tailEnd 在下方；flipV 后 headEnd 在下方。“向下”= 下方那一端带箭头。
    """
    xml = shape_xml(shape)
    flip_v = 'flipV="1"' in xml
    m = re.search(r'prst="(\w*[Aa]rrow\w*)"', xml)
    if m:
        prst = m.group(1).lower()
        # 双向 / 四向 / 折线 / 曲线箭头：方向不稳定，不视作“单向向下”。
        if "leftright" in prst or "quad" in prst or "bent" in prst or "curved" in prst or "updown" in prst:
            return False
        if "down" in prst:
            return not flip_v
        if "up" in prst:
            return flip_v
        return False
    head, tail = arrow_end_types(shape)
    # 竖向单向：仅一端带箭头，且带箭头的那一端在下方。
    if (head is None) == (tail is None):
        return False
    return (tail is not None and not flip_v) or (head is not None and flip_v)


def arrow_points_up(shape) -> bool:
    """判断箭头对象的实际指向是否“向上”。与 arrow_points_down 对称。"""
    xml = shape_xml(shape)
    flip_v = 'flipV="1"' in xml
    m = re.search(r'prst="(\w*[Aa]rrow\w*)"', xml)
    if m:
        prst = m.group(1).lower()
        if "leftright" in prst or "quad" in prst or "bent" in prst or "curved" in prst or "updown" in prst:
            return False
        if "up" in prst:
            return not flip_v
        if "down" in prst:
            return flip_v
        return False
    head, tail = arrow_end_types(shape)
    if (head is None) == (tail is None):
        return False
    # 竖向：未 flipV 时 tailEnd 在下、headEnd 在上；向上 = 上方端带箭头。
    return (head is not None and not flip_v) or (tail is not None and flip_v)


def arrow_points_left(shape) -> bool:
    """判断箭头对象的实际指向是否“向左”。与 is_single_lr_arrow 对称。"""
    xml = shape_xml(shape)
    flip_h = 'flipH="1"' in xml
    m = re.search(r'prst="(\w*[Aa]rrow\w*)"', xml)
    if m:
        prst = m.group(1).lower()
        if "leftright" in prst or "quad" in prst or "bent" in prst or "curved" in prst or "updown" in prst:
            return False
        if "left" in prst:
            return not flip_h
        if "right" in prst:
            return flip_h
        return False
    head, tail = arrow_end_types(shape)
    if (head is None) == (tail is None):
        return False
    return (head is not None and not flip_h) or (tail is not None and flip_h)


def has_dash(shape) -> bool:
    return "prstDash" in shape_xml(shape)


def segment_endpoints(shape) -> tuple[tuple[float, float], tuple[float, float]]:
    """返回线段/连接线/一体箭头的 (起点, 终点)，考虑 flipH/flipV。

    line/cxnSp 默认从 bbox (x1,y1) 画到 (x2,y2)；flipH 交换 x、flipV 交换 y。
    对 bentConnector 之类的一体折线，两个逻辑端点也位于 bbox 对角，同样适用。
    办公软件里“真实呈现的走向”与 flip 状态一致，因此这样得到的端点才是用户看见的起终点。
    """
    x1, y1, x2, y2 = bbox(shape)
    xml = shape_xml(shape)
    flip_h = 'flipH="1"' in xml
    flip_v = 'flipV="1"' in xml
    sx = x2 if flip_h else x1
    sy = y2 if flip_v else y1
    ex = x1 if flip_h else x2
    ey = y1 if flip_v else y2
    return (sx, sy), (ex, ey)


def has_shadow(shape) -> bool:
    xml = shape_xml(shape)
    # 纯粹自动判定阴影较难：有些文件通过主题样式 effectRef 间接应用阴影，
    # XML 中不一定出现 outerShdw。这里把显式阴影和非零 effectRef 都视为满足“带阴影效果”的自动化证据。
    return any(token in xml for token in ["outerShdw", "innerShdw", "effectLst", "effectDag"]) or bool(re.search(r"<a:effectRef\s+idx=\"[1-9]", xml))


def bbox(shape) -> tuple[float, float, float, float]:
    x1, y1 = cm(shape.left), cm(shape.top)
    return x1, y1, x1 + cm(shape.width), y1 + cm(shape.height)


def center(shape) -> tuple[float, float]:
    x1, y1, x2, y2 = bbox(shape)
    return (x1 + x2) / 2, (y1 + y2) / 2


def in_range(v: float, lo: float, hi: float, tol: float = 0.0) -> bool:
    return lo - tol <= v <= hi + tol


def shape_area(shape) -> float:
    return max(cm(shape.width), 0.0) * max(cm(shape.height), 0.0)


def overlap_area(a, b) -> float:
    ax1, ay1, ax2, ay2 = bbox(a)
    bx1, by1, bx2, by2 = bbox(b)
    return max(0, min(ax2, bx2) - max(ax1, bx1)) * max(0, min(ay2, by2) - max(ay1, by1))


def find_shape_by_text(slide, label: str):
    target = norm_text(label)
    best = None
    for sh in iter_shapes(slide):
        t = norm_text(text_of(sh))
        if target and target in t:
            best = sh
            break
    return best


def slide_text(slide) -> str:
    return "\n".join(text_of(s) for s in iter_shapes(slide) if text_of(s))


def find_flow_slide(prs: Presentation) -> tuple[Optional[int], list[int]]:
    candidates = []
    best_idx = None
    best_count = -1
    for idx, slide in enumerate(prs.slides):
        t = norm_text(slide_text(slide))
        count = sum(1 for label in [TARGET_TITLE] + ALL_CORE_LABELS if norm_text(label) in t)
        if count >= 6 or norm_text(TARGET_TITLE) in t:
            candidates.append(idx)
        if count > best_count:
            best_count = count
            best_idx = idx
    return best_idx if best_count >= 6 else None, candidates


def get_lines(slide) -> list:
    # “线条/箭头对象”统一集合：直线连接符、普通线条，以及一体箭头形状。
    # 一体箭头也可作为流程箭头（在办公软件里就是一支完整可选中的箭头），
    # 应与 line/cxnSp 一起纳入候选，以便后续按颜色/长度/方向/端点做匹配。
    return [s for s in iter_shapes(slide) if is_arrow_object(s)]


def get_non_background_shapes(slide, prs) -> list:
    page_area = cm(prs.slide_width) * cm(prs.slide_height)
    result = []
    for s in iter_shapes(slide):
        # 去掉铺底背景框/背景图，避免影响整体流程边界计算。
        if shape_area(s) >= page_area * 0.70 and not text_of(s).strip():
            continue
        if cm(s.width) < 0.05 and cm(s.height) < 0.05 and not text_of(s).strip():
            continue
        result.append(s)
    return result


def find_nearby_text(slide, text: str, x_range=None, y_range=None, font_pt_range=None, color_pred: Callable | None = None) -> Optional[object]:
    target = norm_text(text)
    for sh in iter_shapes(slide):
        if target not in norm_text(text_of(sh)):
            continue
        cx, cy = center(sh)
        if x_range and not in_range(cx, x_range[0], x_range[1]):
            continue
        if y_range and not in_range(cy, y_range[0], y_range[1]):
            continue
        if font_pt_range:
            sizes = text_font_sizes(sh)
            if not sizes or not all(in_range(s, font_pt_range[0], font_pt_range[1], 0.3) for s in sizes):
                continue
        if color_pred:
            colors = text_colors(sh)
            if not colors or not all(color_pred(c) for c in colors):
                continue
        return sh
    return None


def count_icons_near_box(slide, box, allowed_symbols: str = "") -> int:
    """检测框内/附近是否存在可编辑图标。

    评分细则要求图标为可编辑对象。这里用可自动化的替代规则：
    - 独立的小型形状/组合形状，位于流程框内部或紧邻上方；或
    - 文本中含有常见图标符号/emoji（也是可编辑文本）；或
    - 形状名称/XML 含 icon。
    不把图片作为图标计分。
    """
    x1, y1, x2, y2 = bbox(box)
    count = 0
    for sh in iter_shapes(slide):
        if sh is box or is_picture(sh) or is_line(sh):
            continue
        sx, sy = center(sh)
        small = cm(sh.width) <= 1.2 and cm(sh.height) <= 1.2
        inside_or_near = (x1 - 0.15 <= sx <= x2 + 0.15 and y1 - 0.35 <= sy <= y2 + 0.15)
        txt = text_of(sh)
        if small and inside_or_near and not has_chinese(txt):
            count += 1
            continue
        if any(sym in txt for sym in allowed_symbols):
            count += 1
            continue
        if "icon" in (getattr(sh, "name", "") + shape_xml(sh)).lower() and inside_or_near:
            count += 1
    # 有些文件把图标和文字放在同一个文本框中，用 emoji/symbol 作为图标。
    if any(sym in text_of(box) for sym in allowed_symbols):
        count += 1
    return count


def line_length_cm(line) -> float:
    # 一体箭头形状（rightArrow 等）的“长度”应取其主方向的尺寸，而非包围盒最大边——
    # 但对预设箭头图形，主方向恰为其较长的一边（右向箭头 width 更大、上向箭头 height 更大），
    # 因此 max(width, height) 与主方向长度一致，line/cxnSp 亦是如此。
    return max(abs(cm(line.width)), abs(cm(line.height)))


def line_is_horizontal(line, max_dy: float = 0.25) -> bool:
    # 直线连接符/普通线条：其“高度”即为纵向跨度，接近 0 时判为水平。
    # 一体箭头形状：包围盒高度是箭头厚度（如 0.4-1.0cm），因此不能用固定阈值，
    # 改判“宽度显著大于高度”（宽度 > 高度*1.3），即主方向为水平。
    w = abs(cm(line.width))
    h = abs(cm(line.height))
    if is_arrow_shape(line) and not is_line(line):
        return w > 0.2 and w > h * 1.3
    return h <= max_dy and w > 0.05


def line_is_vertical(line, max_dx: float = 0.25) -> bool:
    # 与 line_is_horizontal 对称：一体箭头以“高度显著大于宽度”作为竖向判据。
    w = abs(cm(line.width))
    h = abs(cm(line.height))
    if is_arrow_shape(line) and not is_line(line):
        return h > 0.2 and h > w * 1.3
    return w <= max_dx and h > 0.05


def line_between(line, src, dst, horizontal=True, vertical=False, tol: float = 1.0) -> bool:
    sx, sy = center(src)
    dx, dy = center(dst)
    lx1, ly1, lx2, ly2 = bbox(line)
    lcx, lcy = center(line)
    if horizontal:
        # 箭头通常连接在框的左右边缘，y 坐标应靠近两框垂直中心；允许一定偏差以兼容轻微斜线/连接点偏移。
        return line_is_horizontal(line) and min(sx, dx) - tol <= lcx <= max(sx, dx) + tol and abs(lcy - sy) <= tol * 1.8 and abs(lcy - dy) <= tol * 1.8
    if vertical:
        return line_is_vertical(line) and min(sy, dy) - tol <= lcy <= max(sy, dy) + tol and abs(lcx - sx) <= tol * 1.8 and abs(lcx - dx) <= tol * 1.8
    return min(sx, dx) - tol <= lcx <= max(sx, dx) + tol and min(sy, dy) - tol <= lcy <= max(sy, dy) + tol


def find_horizontal_arrow_between(slide, src, dst, color_pred=is_blue, length_range: tuple[float, float] | None = None, tol=1.0, require_single_lr=False) -> Optional[object]:
    for line in get_lines(slide):
        if not line_between(line, src, dst, horizontal=True, tol=tol):
            continue
        # 颜色判断：一体箭头（填充蓝、无描边色）用 line_rgb 会取到空值被误判失败，
        # 因此对箭头对象统一走 arrow_rgb——它会先取填充、再回退到描边。
        if not color_pred(arrow_rgb(line)):
            continue
        if require_single_lr:
            if not is_single_lr_arrow(line):
                continue
        elif not has_arrowhead(line):
            continue
        if length_range and not in_range(line_length_cm(line), length_range[0], length_range[1], 0.35):
            continue
        return line
    return None


def collect_line_segments(slide, color_pred=is_blue, dashed: Optional[bool] = None) -> list:
    lines = []
    for line in get_lines(slide):
        if not color_pred(line_rgb(line)):
            continue
        if dashed is not None and has_dash(line) != dashed:
            continue
        lines.append(line)
    return lines


@dataclass
class PointResult:
    name: str
    score: int
    ok: bool
    evidence: str


class Evaluator:
    def __init__(self, path: str):
        self.path = path
        self.prs: Optional[Presentation] = None
        self.flow_slide_index: Optional[int] = None
        self.flow_candidates: list[int] = []
        self.gate_results: list[PointResult] = []
        self.points: list[PointResult] = []

    def load(self) -> bool:
        ext = os.path.splitext(self.path)[1].lower()
        if ext != ".pptx":
            # 仅识别 .pptx（Open XML）；二进制 .ppt 需先另存为 .pptx 再评估。
            self.gate_results.append(PointResult("维度1-文件格式与可打开", 0, False, f"扩展名为 {ext}，仅支持 .pptx"))
            return False
        if not os.path.exists(self.path):
            self.gate_results.append(PointResult("维度1-文件格式与可打开", 0, False, "文件不存在"))
            return False
        try:
            if not zipfile.is_zipfile(self.path):
                self.gate_results.append(PointResult("维度1-文件格式与可打开", 0, False, ".pptx 不是有效 ZIP/OpenXML 文件"))
                return False
            self.prs = Presentation(self.path)
            self.gate_results.append(PointResult("维度1-文件格式与可打开", 0, True, f"文件可被 python-pptx 打开，共 {len(self.prs.slides)} 页"))
            return True
        except Exception as e:
            self.gate_results.append(PointResult("维度1-文件格式与可打开", 0, False, f"打开失败：{e}"))
            return False

    @property
    def slide(self):
        assert self.prs is not None and self.flow_slide_index is not None
        return self.prs.slides[self.flow_slide_index]

    def gate_check(self) -> bool:
        if self.prs is None:
            return False
        idx, candidates = find_flow_slide(self.prs)
        self.flow_slide_index = idx
        self.flow_candidates = candidates
        if idx is None:
            self.gate_results.append(PointResult("维度1-包含流程图承载页", 0, False, "未找到包含标题/核心流程文本的幻灯片"))
            return False
        one_flow = len(candidates) == 1
        self.gate_results.append(PointResult(
            "维度1-包含1页流程图承载页",
            0,
            one_flow,
            f"候选承载页：{[i + 1 for i in candidates]}；选用第 {idx + 1} 页" if candidates else "无候选承载页",
        ))
        if not one_flow:
            return False

        return True

    def add_point(self, name: str, score: int, ok: bool, evidence: str) -> None:
        self.points.append(PointResult(name, score, ok, evidence))

    def score_dimension2(self) -> None:
        self.check_title_and_decoration()
        self.check_top_first_three_boxes()
        self.check_top_decision_box()
        self.check_top_last_three_boxes()
        self.check_top_main_arrows()
        self.check_no_label()
        self.check_bottom_four_boxes()
        self.check_bottom_plan_box()
        self.check_bottom_final_decision()
        self.check_bottom_text_format()
        self.check_bottom_arrows()
        self.check_major_down_branch()
        self.check_archive_return_line()
        self.check_archive_return_label()
        self.check_rework_return_line()
        self.check_rework_return_label()
        self.check_record_box()
        self.check_record_text_format()
        self.check_line_styles()
        self.check_editability_score_point()

    def check_title_and_decoration(self):
        slide = self.slide
        sw, shh = cm(self.prs.slide_width), cm(self.prs.slide_height)
        title = find_shape_by_text(slide, TARGET_TITLE)
        title_ok = False
        title_msg = "未找到标题"
        if title is not None:
            cx, cy = center(title)
            sizes = text_font_sizes(title)
            colors = text_colors(title)
            bolds = [r.font.bold for r in all_runs(title) if r.text.strip()]
            title_ok = (
                in_range(cy / shh, 0.05, 0.30, 0.03)                         # 页面上方5%-30%左右
                and abs(cx - sw / 2) <= sw * 0.08                            # 位于页面中间
                and bool(sizes) and all(in_range(s, 32, 34, 0.5) for s in sizes)  # 字号32-34磅
                and bool(bolds) and all(b is True for b in bolds)            # 加粗
                and bool(colors) and all(is_dark_blue(c) for c in colors)    # 深蓝色
                and paragraph_align_center(title)                            # 水平居中
            )
            title_msg = f"中心=({cx:.2f}cm,{cy:.2f}cm={cy/shh:.0%}高)，字号={sizes}，加粗={bolds}，颜色={[rgb_hex(c) for c in colors]}，居中={paragraph_align_center(title)}"

        # 颜色区分：深蓝线/点以蓝色通道为主(b>g)，青绿线/点以绿色通道为主(g>=b)。
        # is_dark_blue 与 is_teal_or_green 判据存在重叠，仅靠色系无法区分左右两条线，
        # 因此叠加“通道主导”判定，对应办公软件里肉眼可见的偏蓝/偏绿。
        def blue_dominant(c):
            return c is not None and c[2] > c[1]

        def green_dominant(c):
            return c is not None and c[1] >= c[2]

        lines = get_lines(slide)
        deco_lines = [ln for ln in lines if in_range(line_length_cm(ln), 4.0, 5.5, 0.05) and line_is_horizontal(ln, 0.1)]
        blue_lines = [ln for ln in deco_lines if is_dark_blue(line_rgb(ln)) and blue_dominant(line_rgb(ln))]
        teal_lines = [ln for ln in deco_lines if is_teal_or_green(line_rgb(ln)) and green_dominant(line_rgb(ln))]

        # 4个小圆点：宽高0.1-0.3cm，按 x 排序后颜色依次为 深蓝、浅蓝、橙色、青绿色。
        dots = [s for s in iter_shapes(slide) if is_oval(s) and in_range(cm(s.width), 0.1, 0.3, 0.01) and in_range(cm(s.height), 0.1, 0.3, 0.01)]
        dots.sort(key=lambda s: center(s)[0])
        dot_colors = [fill_rgb(d) or line_rgb(d) for d in dots]
        dot_seq_ok = (
            len(dots) == 4
            and is_dark_blue(dot_colors[0]) and blue_dominant(dot_colors[0]) and max(dot_colors[0]) <= 140       # 深蓝(暗)
            and is_blue(dot_colors[1]) and blue_dominant(dot_colors[1]) and max(dot_colors[1]) >= 150            # 浅蓝(亮)
            and is_orange(dot_colors[2])                                                                          # 橙色
            and is_teal_or_green(dot_colors[3]) and green_dominant(dot_colors[3])                                 # 青绿色
        )

        # 位置关系：深蓝线在左、青绿线在右；4个圆点位于两线中间；装饰行位于标题下方。
        pos_ok = False
        mid_ok = False
        below_ok = False
        if blue_lines and teal_lines:
            bx = center(blue_lines[0])[0]
            tx = center(teal_lines[0])[0]
            pos_ok = bx < tx
            if dots:
                dcx = sum(center(d)[0] for d in dots) / len(dots)
                mid_ok = min(bx, tx) < dcx < max(bx, tx)
            if title is not None:
                ty = center(title)[1]
                below_ok = all(center(ln)[1] > ty for ln in blue_lines + teal_lines)

        deco_ok = bool(blue_lines) and bool(teal_lines) and dot_seq_ok and pos_ok and mid_ok and below_ok
        self.add_point(
            "+3 标题文本框与标题下装饰线/圆点",
            3,
            title_ok and deco_ok,
            f"标题：{title_msg}；装饰线(4-5.5cm)：深蓝{len(blue_lines)}条、青绿{len(teal_lines)}条，蓝左青右={pos_ok}；"
            f"圆点{len(dots)}个,颜色={[rgb_hex(c) for c in dot_colors]},顺序OK={dot_seq_ok},居中={mid_ok},位标题下方={below_ok}",
        )

    def common_box_text_ok(self, shape, size_lo=10, size_hi=12, center_align=True) -> bool:
        sizes = text_font_sizes(shape)
        colors = text_colors(shape)
        align_ok = paragraph_align_center(shape) if center_align else paragraph_align_left(shape)
        return bool(sizes) and all(in_range(s, size_lo, size_hi, 0.5) for s in sizes) and bool(colors) and all(is_black_or_dark_gray(c) for c in colors) and align_ok and vertical_middle(shape)

    def icon_inside_box(self, slide, box, expected_symbols: str) -> bool:
        """判断某个流程框“框内”是否存在与细则相符、且在办公软件里可编辑的图标。

        细则为每个框指定了具体图标（耳机/文档/柱状图等）。为在 PowerPoint/WPS 中
        真实有效，这里把“图标”限定为真正位于框内的可编辑对象：
        - 一个独立的小型矢量形状（非本框、非图片、非连接线、不含中文文本），其中心
          落在框的范围内（图标叠加在框上）；或
        - 该框自身文本里含有指定类别的图标符号/emoji（也是可编辑文本）。

        注意：不再统计“页面上任意位置存在某符号”——那样一个远处的图标会被所有框
        重复计入，与“框内有X图标”不符，在办公软件里也不成立。
        """
        bx1, by1, bx2, by2 = bbox(box)
        for sh in iter_shapes(slide):
            if sh is box or is_picture(sh) or is_line(sh):
                continue
            sx, sy = center(sh)
            inside = bx1 <= sx <= bx2 and by1 <= sy <= by2
            small = cm(sh.width) <= cm(box.width) * 0.9 and cm(sh.height) <= cm(box.height) * 0.9
            txt = text_of(sh)
            if inside and small and not has_chinese(txt):
                # 独立矢量图标叠加在框内。
                if not txt.strip() or any(sym in txt for sym in expected_symbols):
                    return True
        # 图标以 emoji/符号形式与文字同处该框文本中。
        if any(sym in text_of(box) for sym in expected_symbols):
            return True
        return False

    def check_top_first_three_boxes(self):
        slide = self.slide
        sw = cm(self.prs.slide_width)
        details = []
        ok = True
        # 细则为每个框指定了具体图标：接收投诉=耳机，信息登记=文档，问题分级=柱状图。
        icon_map = {
            "接收投诉": "🎧☎📞",       # 耳机图标
            "信息登记": "📄📋📃🗎",     # 文档图标
            "问题分级": "📊📈",         # 柱状图图标
        }
        for label in TOP_BLUE_LABELS:
            box = find_shape_by_text(slide, label)
            if box is None:
                details.append(f"{label}: 缺失")
                ok = False
                continue
            x1, y1, x2, y2 = bbox(box)
            icon_ok = self.icon_inside_box(slide, box, icon_map[label])
            sizes = text_font_sizes(box)
            colors = text_colors(box)
            size_ok = bool(sizes) and all(in_range(s, 15, 20, 0.5) for s in sizes)         # 字号15-20磅
            color_ok = bool(colors) and all(is_black_or_dark_gray(c) for c in colors)      # 黑色或深灰色
            text_ok = size_ok and color_ok
            box_ok = (
                is_rounded_rect(box)                                 # 形状为圆角矩形
                and in_range(x1 / sw, 0.00, 0.50, 0.01)              # 距左0%-50%
                and in_range(cm(box.width), 3.0, 5.0, 0.05)          # 宽度3-5cm
                and in_range(cm(box.height), 2.0, 3.0, 0.05)         # 高度2-3cm
                and is_dark_blue(line_rgb(box))                      # 轮廓深蓝色
                and has_shadow(box)                                  # 带阴影效果
                and text_ok                                          # 文本15-20磅、黑/深灰
                and icon_ok                                          # 框内有对应图标
            )
            ok = ok and box_ok
            details.append(f"{label}: x={x1/sw:.0%}, 尺寸={cm(box.width):.2f}x{cm(box.height):.2f}cm, 圆角={is_rounded_rect(box)}, 边线={rgb_hex(line_rgb(box))}({'深蓝' if is_dark_blue(line_rgb(box)) else '非深蓝'}), 阴影={has_shadow(box)}, 字号={sizes}({'OK' if size_ok else 'X'}), 颜色={[rgb_hex(c) for c in colors]}({'OK' if color_ok else 'X'}), 指定图标={icon_ok}")
        self.add_point("+5 上方第1-3个圆角矩形流程框", 5, ok, "；".join(details))

    def check_top_decision_box(self):
        slide = self.slide
        sw = cm(self.prs.slide_width)
        box = find_shape_by_text(slide, "是否重大问题？")
        third_box = find_shape_by_text(slide, "问题分级")  # 上方第3个流程框
        ok = False
        msg = "缺失"
        if box is not None:
            x1, y1, x2, y2 = bbox(box)
            mid_y = (y1 + y2) / 2
            # “位于上方第3个流程框右侧”：判断框整体在第3个框右侧。
            right_of_third = third_box is not None and center(box)[0] > center(third_box)[0] and x1 >= bbox(third_box)[2] - 0.3
            # “内部上方有问号图标”：一个独立可编辑对象（非本框），中心落在框内且位于框的上半部，
            # 含问号符号或为小图标形状。区别于框自身的“是否重大问题？”标签文本。
            icon_ok = False
            for sh in iter_shapes(slide):
                if sh is box or is_picture(sh) or is_line(sh):
                    continue
                sx, sy = center(sh)
                inside = x1 <= sx <= x2 and y1 <= sy <= y2
                upper = sy < mid_y
                if not (inside and upper):
                    continue
                t = text_of(sh)
                small = cm(sh.width) <= cm(box.width) * 0.9 and cm(sh.height) <= cm(box.height) * 0.9
                if ("?" in t or "？" in t) or (small and not has_chinese(t)):
                    icon_ok = True
                    break
            sizes = text_font_sizes(box)
            colors = text_colors(box)
            size_ok = bool(sizes) and all(in_range(s, 15, 20, 0.5) for s in sizes)         # 字号15-20磅
            color_ok = bool(colors) and all(is_black_or_dark_gray(c) for c in colors)      # 黑色或深灰色
            text_ok = size_ok and color_ok
            ok = (
                is_diamond(box)                                      # 形状为菱形
                and right_of_third                                   # 位于第3个流程框右侧
                and in_range(x1 / sw, 0.40, 0.60, 0.02)              # 距页面左40%-60%
                and in_range(cm(box.width), 3.0, 4.5, 0.05)          # 宽度3-4.5cm
                and in_range(cm(box.height), 3.0, 4.5, 0.05)         # 高度3-4.5cm
                and is_teal_or_green(line_rgb(box))                  # 轮廓青色或绿色
                and icon_ok                                          # 内部上方有问号图标
                and text_ok                                          # 文本15-20磅、黑/深灰
            )
            msg = f"x={x1/sw:.0%}(需40-60%), 在第3框右侧={right_of_third}, 尺寸={cm(box.width):.2f}x{cm(box.height):.2f}cm, 菱形={is_diamond(box)}, 边线={rgb_hex(line_rgb(box))}({'青/绿' if is_teal_or_green(line_rgb(box)) else '非青绿'}), 上方问号图标={icon_ok}, 字号={sizes}({'OK' if size_ok else 'X'}), 颜色={[rgb_hex(c) for c in colors]}({'OK' if color_ok else 'X'})"
        self.add_point("+1 中部“是否重大问题？”判断框", 1, ok, msg)

    def check_top_last_three_boxes(self):
        slide = self.slide
        sw = cm(self.prs.slide_width)
        decision = find_shape_by_text(slide, "是否重大问题？")  # 中部判断框
        details = []
        ok = True
        # 细则为每个框指定了具体图标：客服处理=人物，客户确认=双人，结案归档=文件夹。
        icon_map = {
            "客服处理": "👤🧑👨👩🎧",     # 人物图标
            "客户确认": "👥👫👬👭🤝",     # 双人图标
            "结案归档": "📁🗂🗃📂",       # 文件夹图标
        }
        for label in ["客服处理", "客户确认", "结案归档"]:
            box = find_shape_by_text(slide, label)
            if box is None:
                details.append(f"{label}: 缺失")
                ok = False
                continue
            x1, y1, x2, y2 = bbox(box)
            # “位于判断框右侧”：框整体在中部判断框右侧。
            right_of_decision = decision is not None and center(box)[0] > center(decision)[0] and x1 >= bbox(decision)[2] - 0.3
            icon_ok = self.icon_inside_box(slide, box, icon_map[label])
            sizes = text_font_sizes(box)
            colors = text_colors(box)
            size_ok = bool(sizes) and all(in_range(s, 15, 20, 0.5) for s in sizes)         # 字号15-20磅
            color_ok = bool(colors) and all(is_black_or_dark_gray(c) for c in colors)      # 黑色或深灰色
            text_ok = size_ok and color_ok
            box_ok = (
                is_rounded_rect(box)                                 # 圆角矩形
                and right_of_decision                                # 位于判断框右侧
                and in_range(x1 / sw, 0.55, 1.00, 0.02)              # 距页面左55%-100%
                and in_range(cm(box.width), 3.0, 5.0, 0.05)          # 宽度3-5cm
                and in_range(cm(box.height), 2.0, 3.0, 0.05)         # 高度2-3cm
                and is_teal_or_green(line_rgb(box))                  # 轮廓青色或绿色
                and has_shadow(box)                                  # 带阴影效果
                and text_ok                                          # 文本15-20磅、黑/深灰
                and icon_ok                                          # 框内有对应图标
            )
            ok = ok and box_ok
            details.append(f"{label}: x={x1/sw:.0%}, 判断框右侧={right_of_decision}, 尺寸={cm(box.width):.2f}x{cm(box.height):.2f}cm, 圆角={is_rounded_rect(box)}, 边线={rgb_hex(line_rgb(box))}({'青/绿' if is_teal_or_green(line_rgb(box)) else '非青绿'}), 阴影={has_shadow(box)}, 字号={sizes}({'OK' if size_ok else 'X'}), 颜色={[rgb_hex(c) for c in colors]}({'OK' if color_ok else 'X'}), 指定图标={icon_ok}")
        self.add_point("+5 上方第4-6个圆角矩形流程框", 5, ok, "；".join(details))

    def check_top_main_arrows(self):
        slide = self.slide
        # 细则逐段规定了长度；均要求蓝色、单向、水平、方向从左向右。
        pairs = [
            ("接收投诉", "信息登记", (0.5, 1.3)),
            ("信息登记", "问题分级", (0.5, 1.3)),
            ("问题分级", "是否重大问题？", (0.5, 1.0)),
            ("是否重大问题？", "客服处理", (1.1, 1.7)),
            ("客服处理", "客户确认", (0.5, 1.3)),
            ("客户确认", "结案归档", (0.5, 1.3)),
        ]
        ok = True
        details = []
        for a, b, lr in pairs:
            src, dst = find_shape_by_text(slide, a), find_shape_by_text(slide, b)
            # require_single_lr：不仅要有箭头，且必须是“单向、从左向右”的水平箭头。
            ln = find_horizontal_arrow_between(slide, src, dst, is_blue, lr, tol=1.2, require_single_lr=True) if src and dst else None
            ok = ok and ln is not None
            details.append(f"{a}->{b}: {'命中' if ln else '未命中'}" + (f"，长度={line_length_cm(ln):.2f}cm，颜色={rgb_hex(line_rgb(ln))}，单向右向={is_single_lr_arrow(ln)}" if ln else ""))
        self.add_point("+3 上方主流程蓝色水平单向箭头", 3, ok, "；".join(details))

    def check_no_label(self):
        slide = self.slide
        # 细则：在“判断框通向客服处理的箭头”上方或附近出现“否”，字号14-16磅、黑/深灰。
        # 定位应以那条箭头为基准，而不是写死坐标——这样才对应办公软件里的真实相对位置。
        decision = find_shape_by_text(slide, "是否重大问题？")
        kf = find_shape_by_text(slide, "客服处理")
        arrow = find_horizontal_arrow_between(slide, decision, kf, is_blue, (1.1, 1.7), tol=1.2, require_single_lr=True) if decision and kf else None

        ok = False
        msg = "未找到判断框→客服处理的箭头，无法定位“否”"
        if arrow is not None:
            ax, ay = center(arrow)
            found = None
            for sh in iter_shapes(slide):
                t = norm_text(text_of(sh))
                if t != "否":
                    continue
                cx, cy = center(sh)
                # “上方或附近”：水平靠近箭头，且位于箭头上方或紧邻（不在其显著下方）。
                near = abs(cx - ax) <= 1.5 and (cy <= ay + 0.3) and (ay - cy) <= 2.5
                if not near:
                    continue
                sizes = text_font_sizes(sh)
                colors = text_colors(sh)
                size_ok = bool(sizes) and all(in_range(s, 15, 25, 0.5) for s in sizes)   # 字号15-25磅
                color_ok = bool(colors) and all(is_black_or_dark_gray(c) for c in colors)  # 黑/深灰
                if size_ok and color_ok:
                    found = sh
                    break
                msg = f"找到“否”于({cx:.2f},{cy:.2f})，字号={sizes}(需15-25)，颜色={[rgb_hex(c) for c in colors]}"
            ok = found is not None
            if ok:
                msg = f"“否”位于箭头({ax:.2f},{ay:.2f})上方/附近，字号={text_font_sizes(found)}，颜色={[rgb_hex(c) for c in text_colors(found)]}"
            elif msg.startswith("未找到判断框"):
                msg = f"箭头({ax:.2f},{ay:.2f})上方/附近未找到符合格式的“否”"
        self.add_point("+1 判断框右上分支标注“否”", 1, ok, msg)

    def check_bottom_four_boxes(self):
        slide = self.slide
        shh = cm(self.prs.slide_height)
        details = []
        ok = True
        # 细则为每个框指定了具体图标：指派专员=人物，原因分析=放大镜，执行整改=齿轮，质量复核=盾牌对勾。
        icon_map = {
            "指派专员": "👤🧑👨👩",       # 人物图标
            "原因分析": "🔍🔎",           # 放大镜图标
            "执行整改": "⚙🛠🔧",          # 齿轮图标
            "质量复核": "🛡✅✔☑",         # 盾牌对勾图标
        }
        for label in ["指派专员", "原因分析", "执行整改", "质量复核"]:
            box = find_shape_by_text(slide, label)
            if box is None:
                details.append(f"{label}: 缺失")
                ok = False
                continue
            x1, y1, x2, y2 = bbox(box)
            icon_ok = self.icon_inside_box(slide, box, icon_map[label])
            box_ok = (
                is_rounded_rect(box)                                 # 形状为圆角矩形
                and in_range(y1 / shh, 0.60, 0.80, 0.02)             # 距上60%-80%
                and in_range(cm(box.width), 2.5, 3.5, 0.08)          # 宽度2.5-3.5cm
                and in_range(cm(box.height), 2.5, 3.5, 0.08)         # 高度2.5-3.5cm
                and is_orange(line_rgb(box))                         # 轮廓橙色
                and has_shadow(box)                                  # 带阴影效果
                and icon_ok                                          # 框内有对应图标
            )
            ok = ok and box_ok
            details.append(f"{label}: y={y1/shh:.0%}, 尺寸={cm(box.width):.2f}x{cm(box.height):.2f}cm, 圆角={is_rounded_rect(box)}, 边线={rgb_hex(line_rgb(box))}({'橙' if is_orange(line_rgb(box)) else '非橙'}), 阴影={has_shadow(box)}, 指定图标={icon_ok}")
        self.add_point("+5 下方第1、2、4、5个整改流程框", 5, ok, "；".join(details))

    def check_bottom_plan_box(self):
        slide = self.slide
        # 细则：位于“原因分析”右侧；宽2.5-3.5cm、高2.5-3.5cm；圆角矩形；橙色轮廓；带阴影效果；
        #       框内有清单图标和文本“制定整改方案”。
        # 通过文本定位该框，即隐含校验“框内有文本‘制定整改方案’”。
        box = find_shape_by_text(slide, "制定整改方案")
        cause = find_shape_by_text(slide, "原因分析")
        ok = False
        msg = "缺失"
        if box is not None:
            x1, y1, x2, y2 = bbox(box)
            # “位于原因分析右侧”：本框整体在“原因分析”框右侧（中心更靠右且左边缘不越过其右边缘），
            # 对应办公软件里两框的真实左右相对位置。
            right_of_cause = cause is not None and center(box)[0] > center(cause)[0] and x1 >= bbox(cause)[2] - 0.3
            # “框内有清单图标”：清单类图标须为真正落在框内的可编辑对象，或本框文本内含清单符号，
            # 而非页面别处的图标——这样在 PowerPoint/WPS 中“框内有图标”才成立。
            icon_ok = self.icon_inside_box(slide, box, "📋☑☒☐✓✔")
            ok = (
                is_rounded_rect(box)                             # 形状为圆角矩形
                and right_of_cause                               # 位于“原因分析”右侧
                and in_range(cm(box.width), 2.5, 3.5, 0.08)      # 宽度2.5-3.5cm
                and in_range(cm(box.height), 2.5, 3.5, 0.08)     # 高度2.5-3.5cm
                and is_orange(line_rgb(box))                     # 轮廓橙色
                and has_shadow(box)                              # 带阴影效果
                and icon_ok                                      # 框内有清单图标
            )
            msg = f"尺寸={cm(box.width):.2f}x{cm(box.height):.2f}cm, 圆角={is_rounded_rect(box)}, 原因分析右侧={right_of_cause}, 边线={rgb_hex(line_rgb(box))}({'橙' if is_orange(line_rgb(box)) else '非橙'}), 阴影={has_shadow(box)}, 框内清单图标={icon_ok}"
        self.add_point("+1 下方第3个“制定整改方案”流程框", 1, ok, msg)

    def check_bottom_final_decision(self):
        slide = self.slide
        # 细则：位于“质量复核”右侧；宽3.4-4.5cm、高3.5-4.5cm；菱形；橙色轮廓；
        #       内部有带勾的清单图标和文本“复核通过？”。
        # 通过文本定位该框，即隐含校验“内部有文本‘复核通过？’”。
        box = find_shape_by_text(slide, "复核通过？")
        review = find_shape_by_text(slide, "质量复核")
        ok = False
        msg = "缺失"
        if box is not None:
            x1, y1, x2, y2 = bbox(box)
            # “位于质量复核右侧”：本框整体在“质量复核”框右侧（中心更靠右且左边缘不越过其右边缘），
            # 对应办公软件里两框的真实左右相对位置。
            right_of_review = review is not None and center(box)[0] > center(review)[0] and x1 >= bbox(review)[2] - 0.3
            # “内部有带勾的清单图标”：带勾清单类图标须为真正落在框内的可编辑对象，或本框文本内含带勾清单符号，
            # 而非页面别处的图标——这样在 PowerPoint/WPS 中“框内有图标”才成立。
            icon_ok = self.icon_inside_box(slide, box, "✅✔✓☑📋")
            ok = (
                is_diamond(box)                                  # 形状为菱形
                and right_of_review                              # 位于“质量复核”右侧
                and in_range(cm(box.width), 3.4, 4.5, 0.08)      # 宽度3.4-4.5cm
                and in_range(cm(box.height), 3.5, 4.5, 0.08)     # 高度3.5-4.5cm
                and is_orange(line_rgb(box))                     # 轮廓橙色
                and icon_ok                                      # 内部有带勾的清单图标
            )
            msg = f"尺寸={cm(box.width):.2f}x{cm(box.height):.2f}cm, 菱形={is_diamond(box)}, 质量复核右侧={right_of_review}, 边线={rgb_hex(line_rgb(box))}({'橙' if is_orange(line_rgb(box)) else '非橙'}), 框内带勾清单图标={icon_ok}"
        self.add_point("+1 下方终点“复核通过？”判断框", 1, ok, msg)

    def check_bottom_text_format(self):
        slide = self.slide
        # 细则：下方5个矩形框(指派专员/原因分析/制定整改方案/执行整改/质量复核)
        #       与1个菱形框(复核通过？)中的文字——字号15-20磅、颜色黑色或深灰色。
        # 逐框、逐 run 校验，且判据均为办公软件(PowerPoint/WPS)里真实生效的属性：
        #   - 字号取 run.font.size（磅）
        #   - 颜色取 run 实际字体颜色
        labels = BOTTOM_LABELS
        details = []
        ok = True
        for label in labels:
            box = find_shape_by_text(slide, label)
            if box is None:
                ok = False
                details.append(f"{label}: 缺失")
                continue
            sizes = text_font_sizes(box)
            colors = text_colors(box)
            size_ok = bool(sizes) and all(in_range(s, 15, 20, 0.5) for s in sizes)   # 字号15-20磅
            color_ok = bool(colors) and all(is_black_or_dark_gray(c) for c in colors)  # 黑色/深灰色
            b_ok = size_ok and color_ok
            ok = ok and b_ok
            details.append(f"{label}: 字号={sizes}({'OK' if size_ok else 'X'}), 颜色={[rgb_hex(c) for c in colors]}({'OK' if color_ok else 'X'})")
        self.add_point("+1 下方整改流程文字格式", 1, ok, "；".join(details))

    def check_bottom_arrows(self):
        slide = self.slide
        # 细则：下列相邻两框之间均有——长度0.5-1.3cm、蓝色、单向、水平、方向从左向右的箭头。
        pairs = [
            ("指派专员", "原因分析"),
            ("原因分析", "制定整改方案"),
            ("制定整改方案", "执行整改"),
            ("执行整改", "质量复核"),
            ("质量复核", "复核通过？"),
        ]
        ok = True
        details = []
        for a, b in pairs:
            src, dst = find_shape_by_text(slide, a), find_shape_by_text(slide, b)
            # require_single_lr：不仅要有箭头，且必须是“单向、从左向右”的水平箭头，
            # 这是办公软件(PowerPoint/WPS)里真实生效的箭头方向。
            ln = find_horizontal_arrow_between(slide, src, dst, is_blue, (0.5, 1.3), tol=1.2, require_single_lr=True) if src and dst else None
            ok = ok and ln is not None
            details.append(f"{a}->{b}: {'命中' if ln else '未命中'}" + (f"，长度={line_length_cm(ln):.2f}cm，颜色={rgb_hex(line_rgb(ln))}，单向右向={is_single_lr_arrow(ln)}" if ln else ""))
        self.add_point("+3 下方整改流程蓝色水平单向箭头", 3, ok, "；".join(details))

    def check_major_down_branch(self):
        slide = self.slide
        # 细则：上方“是否重大问题？”菱形框底部向下连接1条 长度1.7-2.5cm、蓝色、竖向的箭头，
        #       指向“制定整改方案”上方或该下方流程区域；箭头附近出现文本“是”。
        decision = find_shape_by_text(slide, "是否重大问题？")
        plan = find_shape_by_text(slide, "制定整改方案")
        arrow = None
        if decision is not None:
            dx1, dy1, dx2, dy2 = bbox(decision)
            dcx = (dx1 + dx2) / 2
            for ln in get_lines(slide):
                if not line_is_vertical(ln):                          # 竖向
                    continue
                if not arrow_points_down(ln):                         # 方向必须向下（避免向上竖箭头误判通过）
                    continue
                if not is_blue(arrow_rgb(ln)):                        # 蓝色（一体箭头用填充色）
                    continue
                if not has_arrowhead(ln):                             # 箭头
                    continue
                if not in_range(line_length_cm(ln), 1.7, 2.5):        # 长度1.7-2.5cm
                    continue
                lcx, lcy = center(ln)
                lx1, ly1, lx2, ly2 = bbox(ln)
                # 从菱形框底部向下：横向与菱形对齐，且整体位于菱形下方（顶端不高于菱形底边）。
                if abs(lcx - dcx) > 1.5:
                    continue
                if ly1 < dy2 - 0.5:
                    continue
                # 指向“制定整改方案”上方或该下方流程区域：竖箭头横向靠近该框，
                # 末端到达其上方或其所在整改流程行范围内。
                points_to_plan_area = True
                if plan is not None:
                    px1, py1, px2, py2 = bbox(plan)
                    points_to_plan_area = abs(lcx - (px1 + px2) / 2) <= 3.0 and ly2 <= py2 + 0.5
                if not points_to_plan_area:
                    continue
                arrow = ln
                break
        # 箭头附近出现文本“是”：以箭头为基准定位（对应办公软件里的真实相对位置），不写死坐标。
        yes = None
        if arrow is not None:
            ax, ay = center(arrow)
            for sh in iter_shapes(slide):
                if norm_text(text_of(sh)) != "是":
                    continue
                cx, cy = center(sh)
                if abs(cx - ax) <= 2.0 and abs(cy - ay) <= 2.0:
                    yes = sh
                    break
        ok = arrow is not None and yes is not None
        msg = f"菱形下方竖向蓝色箭头(1.7-2.5cm,指向整改区)={arrow is not None}；箭头附近“是”标注={yes is not None}"
        self.add_point("+1 “是否重大问题？”向下“是”分支", 1, ok, msg)

    def check_archive_return_line(self):
        slide = self.slide
        # 细则：下方“复核通过？”菱形框右侧或右上侧，通过蓝色折线箭头连接至上方“结案归档”框下侧；
        #       箭头最终指向“结案归档”。
        #
        # 严格做法：按端点顺序重建折线，验证——
        #   1) 起点位于“复核通过？”右侧或右上侧；
        #   2) 折线由若干段首尾相接组成，方向依次为“向右 → 向上 → 向左”（允许单段折线连接器直接跨越）；
        #   3) 末段带箭头，末端落在“结案归档”底部附近，且末段方向为“向上”或“向左”（对应箭头真正指向归档）。
        review = find_shape_by_text(slide, "复核通过？")
        archive = find_shape_by_text(slide, "结案归档")
        lines = collect_line_segments(slide, is_blue, dashed=False)

        start_ok = False
        arrow_to_archive = False
        path_ok = False
        path_debug = ""

        if review is not None and archive is not None and lines:
            rx1, ry1, rx2, ry2 = bbox(review)
            rcy_mid = (ry1 + ry2) / 2
            ax1, ay1, ax2, ay2 = bbox(archive)
            acx = (ax1 + ax2) / 2

            # 每段线的两个逻辑端点（考虑 flipH/flipV），并按“主方向”分类，供路径重建使用。
            segs = []
            for ln in lines:
                (sx, sy), (ex, ey) = segment_endpoints(ln)
                dx = ex - sx
                dy = ey - sy
                if abs(dx) < 0.15 and abs(dy) < 0.15:
                    continue  # 过短，无意义
                if abs(dx) >= abs(dy):
                    direction = "right" if dx > 0 else "left"
                else:
                    direction = "down" if dy > 0 else "up"
                segs.append({
                    "shape": ln,
                    "start": (sx, sy),
                    "end": (ex, ey),
                    "dir": direction,
                    "arrow": has_arrowhead(ln),
                    "len": max(abs(dx), abs(dy)),
                })

            def near(p, q, tol=0.6) -> bool:
                return abs(p[0] - q[0]) <= tol and abs(p[1] - q[1]) <= tol

            # 起点必须位于“复核通过？”右侧（右缘附近）或右上侧（右缘且不低于纵向中线）。
            def is_valid_start(p) -> bool:
                px, py = p
                on_right = px >= rx1 - 0.3 and px <= rx2 + 1.5
                not_below_mid = py <= rcy_mid + 0.5
                return on_right and not_below_mid

            # 终点必须落在“结案归档”底部附近（下缘中段）。
            def is_valid_end(p) -> bool:
                px, py = p
                return abs(px - acx) <= 2.5 and in_range(py, ay2 - 0.5, ay2 + 2.5)

            # 依次尝试每段作为路径首段——起点合法、方向为“向右”（右侧接出后先水平向外）。
            # 然后按端点连通性做至多 4 段的贪心链接；每一步只允许期望的方向序列。
            # 序列模式一：[right, up, left]（标准回流折线）；
            # 序列模式二：[right, up]（折线只两段，或一体折线连接器分成右+左上的两段）；
            # 序列模式三：[up]（起点直接在右上、单段竖箭头上行到归档下方——较少见但兼容）。
            expected_seqs = [
                ["right", "up", "left"],
                ["right", "up"],
                ["up", "left"],
                ["up"],
            ]

            def try_chain(expected: list) -> Optional[list]:
                # 找到一条端点连通、方向依次匹配 expected、末段末端指向归档底部、末段带箭头的路径。
                for i, s0 in enumerate(segs):
                    if s0["dir"] != expected[0]:
                        continue
                    if not is_valid_start(s0["start"]):
                        continue
                    chain = [s0]
                    cursor = s0["end"]
                    used = {i}
                    ok_chain = True
                    for step in expected[1:]:
                        nxt = None
                        for j, sj in enumerate(segs):
                            if j in used:
                                continue
                            if sj["dir"] != step:
                                continue
                            # 端点连通：下一段起点与当前游标接近，或该段两端任一接近（兼容 flip 未识别的情况）。
                            if near(sj["start"], cursor):
                                nxt = (j, sj, sj["end"])
                                break
                            if near(sj["end"], cursor):
                                # 反过来接入：把该段视为从 end 走向 start。
                                nxt = (j, sj, sj["start"])
                                break
                        if nxt is None:
                            ok_chain = False
                            break
                        j, sj, new_cursor = nxt
                        used.add(j)
                        chain.append(sj)
                        cursor = new_cursor
                    if not ok_chain:
                        continue
                    last = chain[-1]
                    if not last["arrow"]:
                        continue
                    # 末端必须落在归档底部附近——用 cursor（即末段推进后的实际末端）判定。
                    if not is_valid_end(cursor):
                        continue
                    # 末段方向必须朝向归档（向上或向左），避免出现末段方向为“向下/向右”仍误判通过。
                    if last["dir"] not in ("up", "left"):
                        continue
                    return chain
                return None

            chain = None
            for seq in expected_seqs:
                chain = try_chain(seq)
                if chain is not None:
                    path_ok = True
                    start_ok = True
                    arrow_to_archive = True
                    path_debug = "→".join(s["dir"] for s in chain)
                    break

            if chain is None:
                # 用于失败时的诊断输出：把关键的三个判据分别单独列出，便于用户排查。
                start_ok = any(is_valid_start(s["start"]) and s["dir"] == "right" for s in segs) or \
                           any(is_valid_start(s["start"]) and s["dir"] == "up" for s in segs)
                arrow_to_archive = any(s["arrow"] and is_valid_end(s["end"]) and s["dir"] in ("up", "left") for s in segs)

        ok = path_ok
        self.add_point(
            "+1 结案归档回流蓝色折线箭头",
            1,
            ok,
            f"路径重建={path_ok}({path_debug})，起点位于复核通过右侧/右上={start_ok}，末段箭头指向结案归档下侧={arrow_to_archive}，蓝色实线段数={len(lines)}",
        )

    def check_archive_return_label(self):
        slide = self.slide
        # 细则：从“复核通过？”连向“结案归档”的折线附近出现文本“是”，字号14-25磅、颜色黑色或深灰色。
        # 定位以那条回流折线为基准——折线大致位于“复核通过？”与“结案归档”之间的区域，
        # 这样对应办公软件里“折线附近”的真实相对位置，而不写死坐标。
        review = find_shape_by_text(slide, "复核通过？")
        archive = find_shape_by_text(slide, "结案归档")
        yes = None
        msg = "未找到“复核通过？”或“结案归档”，无法定位折线附近的“是”"
        if review is not None and archive is not None:
            rcx, rcy = center(review)
            acx, acy = center(archive)
            # 折线附近区域：横向在两框之间偏右侧（回流线上行处），纵向介于两框之间。
            zone_x_lo = min(rcx, acx) - 0.5
            zone_x_hi = max(rcx, acx) + 2.5
            zone_y_lo = min(rcy, acy) - 1.0
            zone_y_hi = max(rcy, acy) + 1.0
            msg = "折线附近未找到符合字号/颜色要求的“是”"
            for sh in iter_shapes(slide):
                if norm_text(text_of(sh)) != "是":
                    continue
                cx, cy = center(sh)
                if not (zone_x_lo <= cx <= zone_x_hi and zone_y_lo <= cy <= zone_y_hi):
                    continue
                sizes = text_font_sizes(sh)
                colors = text_colors(sh)
                size_ok = bool(sizes) and all(in_range(s, 14, 25, 0.5) for s in sizes)      # 字号14-25磅
                color_ok = bool(colors) and all(is_black_or_dark_gray(c) for c in colors)   # 黑色/深灰色
                if size_ok and color_ok:
                    yes = sh
                    msg = f"折线附近找到“是”于({cx:.2f},{cy:.2f})，字号={sizes}，颜色={[rgb_hex(c) for c in colors]}"
                    break
                msg = f"折线附近找到“是”于({cx:.2f},{cy:.2f})，字号={sizes}(需14-25)，颜色={[rgb_hex(c) for c in colors]}"
        self.add_point("+1 结案归档回流线标注“是”", 1, yes is not None, msg)

    def check_rework_return_line(self):
        slide = self.slide
        # 细则：从“复核通过？”菱形框下方引出1条蓝色虚线折线箭头，
        #       折线 向下2-3cm 后 向左7.5-9cm 延伸，再 向上2-3cm 连接到“执行整改”框底部。
        review = find_shape_by_text(slide, "复核通过？")
        exec_box = find_shape_by_text(slide, "执行整改")
        lines = collect_line_segments(slide, is_blue, dashed=True)   # 蓝色虚线
        down = [ln for ln in lines if line_is_vertical(ln) and in_range(line_length_cm(ln), 2.0, 3.0)]
        left = [ln for ln in lines if line_is_horizontal(ln) and in_range(line_length_cm(ln), 7.5, 9.0)]
        up = [ln for ln in lines if line_is_vertical(ln) and in_range(line_length_cm(ln), 2.0, 3.0)]

        # 起于“复核通过？”下方：存在一条向下段，横向对齐菱形、位于其底边下方。
        start_ok = False
        if review is not None and down:
            rx1, ry1, rx2, ry2 = bbox(review)
            rcx = (rx1 + rx2) / 2
            for ln in down:
                lcx, lcy = center(ln)
                if abs(lcx - rcx) <= 1.5 and lcy >= ry2 - 0.5:
                    start_ok = True
                    break
        # 连接到“执行整改”框底部：存在一条向上段(带箭头)，横向对齐该框、末端到达其底边附近。
        end_ok = False
        if exec_box is not None and up:
            ex1, ey1, ex2, ey2 = bbox(exec_box)
            ecx = (ex1 + ex2) / 2
            for ln in up:
                if not has_arrowhead(ln):
                    continue
                lcx, lcy = center(ln)
                lx1, ly1, lx2, ly2 = bbox(ln)
                if abs(lcx - ecx) <= 1.5 and in_range(ly1, ey2 - 0.5, ey2 + 2.5):
                    end_ok = True
                    break
        ok = start_ok and bool(down) and bool(left) and bool(up) and end_ok
        self.add_point(
            "+1 返工回流蓝色虚线折线箭头",
            1,
            ok,
            f"起于复核通过下方={start_ok}，虚线向下2-3cm段={len(down)}，虚线向左7.5-9cm段={len(left)}，向上2-3cm段={len(up)}，箭头连执行整改底部={end_ok}",
        )

    def check_rework_return_label(self):
        slide = self.slide
        # 细则：
        #  1) 虚线回路线中部或右侧出现文本“返工”，位于“执行整改”下方附近；
        #  2) “复核通过？”菱形框下方附近出现文本“否”，字号14-25磅、颜色黑色或深灰色。
        # 均以相关流程框的真实位置为基准定位，对应办公软件里的相对位置，不写死坐标。
        exec_box = find_shape_by_text(slide, "执行整改")
        review = find_shape_by_text(slide, "复核通过？")

        # “返工”：位于“执行整改”下方附近（横向靠近该框、纵向在其下方）。
        rework = None
        if exec_box is not None:
            ex1, ey1, ex2, ey2 = bbox(exec_box)
            ecx = (ex1 + ex2) / 2
            for sh in iter_shapes(slide):
                if norm_text(text_of(sh)) != "返工":
                    continue
                cx, cy = center(sh)
                if abs(cx - ecx) <= 4.0 and ey2 - 0.5 <= cy <= ey2 + 4.0:
                    rework = sh
                    break

        # “否”：位于“复核通过？”菱形框下方附近，且字号14-25磅、黑/深灰。
        no = None
        if review is not None:
            rx1, ry1, rx2, ry2 = bbox(review)
            rcx = (rx1 + rx2) / 2
            for sh in iter_shapes(slide):
                if norm_text(text_of(sh)) != "否":
                    continue
                cx, cy = center(sh)
                if not (abs(cx - rcx) <= 2.5 and ry2 - 0.5 <= cy <= ry2 + 3.0):
                    continue
                sizes = text_font_sizes(sh)
                colors = text_colors(sh)
                size_ok = bool(sizes) and all(in_range(s, 14, 25, 0.5) for s in sizes)      # 字号14-25磅
                color_ok = bool(colors) and all(is_black_or_dark_gray(c) for c in colors)   # 黑色/深灰色
                if size_ok and color_ok:
                    no = sh
                    break
        ok = rework is not None and no is not None
        self.add_point("+1 返工回路线标注“返工/否”", 1, ok, f"“返工”位于执行整改下方={rework is not None}；“否”位于复核通过下方且格式合规={no is not None}")

    def check_record_box(self):
        slide = self.slide
        sw, shh = cm(self.prs.slide_width), cm(self.prs.slide_height)
        # 细则：
        #  1) 页面左下区域出现1个矩形框；2) 白底；3) 蓝色轮廓；4) 宽5-7cm；5) 高4-5cm；且为竖向（高不小于宽）；
        #  6) 框内顶部文本为“相关记录：”；7) 下方分3行项目符号文本“沟通记录”“证据材料”“处理报告”；
        #  8) 说明框左下外侧或框内左侧出现蓝色圆形图标，圆内为剪贴板图案，且为可编辑形状组合（非图片）。

        # —— 记录说明框：左下区域、矩形、白底、蓝色轮廓、尺寸 5-7 x 4-5 cm、竖向 ——
        box = None
        for sh in iter_shapes(slide):
            if not is_rectangle(sh):
                continue
            x1, y1, x2, y2 = bbox(sh)
            if not (x1 < sw * 0.35 and y1 > shh * 0.55):          # 页面左下区域
                continue
            if not is_white_fill(fill_rgb(sh)):                   # 白底
                continue
            if not is_blue(line_rgb(sh)):                         # 蓝色轮廓
                continue
            if not (in_range(cm(sh.width), 5.0, 7.0) and in_range(cm(sh.height), 4.0, 5.0)):  # 5-7 x 4-5 cm
                continue
            box = sh
            break

        # —— 框内文本：顶部“相关记录：”，其下3行项目符号“沟通记录/证据材料/处理报告” ——
        text_ok = False
        top_ok = False
        bullets_ok = False
        if box is not None:
            bx1, by1, bx2, by2 = bbox(box)
            rec_shape = None
            for sh in iter_shapes(slide):
                if not getattr(sh, "has_text_frame", False):
                    continue
                if norm_text("相关记录：") not in norm_text(text_of(sh)):
                    continue
                scx, scy = center(sh)
                if bx1 - 0.5 <= scx <= bx2 + 0.5 and by1 - 0.5 <= scy <= by2 + 0.5:  # 文本位于框内
                    rec_shape = sh
                    break
            if rec_shape is not None:
                paras = [p for p in rec_shape.text_frame.paragraphs if p.text.strip()]
                if paras:
                    # 顶部第一段为“相关记录：”
                    top_ok = norm_text(paras[0].text).startswith("相关记录")
                    # 其下的项目符号行文本（按出现顺序）应依次含 沟通记录/证据材料/处理报告
                    def is_bullet_para(p) -> bool:
                        try:
                            xml = etree.tostring(p._p, encoding="unicode")
                        except Exception:
                            xml = ""
                        if re.search(r"<a:(buChar|buAutoNum)\b", xml):
                            return True
                        return norm_text(p.text)[:1] in "•‣◦·▪-*"
                    wanted = ["沟通记录", "证据材料", "处理报告"]
                    bullet_paras = [p for p in paras[1:] if is_bullet_para(p)]
                    hit = [w for w in wanted if any(w in norm_text(p.text) for p in bullet_paras)]
                    bullets_ok = len(bullet_paras) >= 3 and len(hit) == 3
            text_ok = top_ok and bullets_ok

        # —— 蓝色圆形剪贴板图标：位置在框左侧内或左下外侧；蓝色；圆形；圆内剪贴板图案；可编辑形状（非图片）——
        icon_ok = False
        icon_msg = "未检测"
        clipboard_syms = "📋🗒🗐📄🖇"
        if box is not None:
            bx1, by1, bx2, by2 = bbox(box)
            box_left_mid = bx1 + (bx2 - bx1) * 0.5
            for sh in iter_shapes(slide):
                if sh is box or is_picture(sh) or is_line(sh):     # 不接受图片（须为可编辑形状）
                    continue
                # 圆形（圆形本身，或包含圆形子形状的可编辑组合）
                circular = is_oval(sh) and abs(cm(sh.width) - cm(sh.height)) <= max(cm(sh.width), 0.01) * 0.35
                group_with_oval = is_group(sh) and any(is_oval(c) for c in iter_shapes(sh))
                if not (circular or group_with_oval):
                    continue
                # 蓝色（填充或轮廓为蓝色系）
                if not (is_blue(fill_rgb(sh)) or is_blue(line_rgb(sh))):
                    continue
                icx, icy = center(sh)
                inside_left = (bx1 - 0.3 <= icx <= box_left_mid) and (by1 - 0.3 <= icy <= by2 + 0.3)   # 框内左侧
                lower_left_outside = (icx <= bx1 + 0.3) and (icy >= by2 - 1.0)                          # 左下外侧
                if not (inside_left or lower_left_outside):
                    continue
                # 圆内为剪贴板图案：本形状文本含剪贴板符号，或其组合内含带剪贴板符号的子形状。
                own_txt = text_of(sh)
                sub_txt = "".join(text_of(c) for c in iter_shapes(sh)) if is_group(sh) else ""
                has_clip = any(sym in own_txt for sym in clipboard_syms) or any(sym in sub_txt for sym in clipboard_syms)
                if not has_clip:
                    continue
                icon_ok = True
                icon_msg = f"蓝色圆形+剪贴板，位置={'框内左侧' if inside_left else '左下外侧'}"
                break
            if not icon_ok:
                icon_msg = "未找到符合(蓝色/圆形/剪贴板/可编辑)的图标"

        ok = box is not None and text_ok and icon_ok
        if box is None:
            msg = "缺失符合(左下/白底/蓝框/5-7x4-5cm)的记录说明框"
        else:
            msg = f"说明框尺寸={cm(box.width):.2f}x{cm(box.height):.2f}cm，白底={is_white_fill(fill_rgb(box))}，边线={rgb_hex(line_rgb(box))}，顶部相关记录={top_ok}，3行项目符号={bullets_ok}，图标：{icon_msg}"
        self.add_point("+5 左下相关记录说明框及剪贴板图标", 5, ok, msg)

    def check_record_text_format(self):
        slide = self.slide
        # 细则：说明框中的“标题”和“项目符号文字”——字号15-22磅、颜色黑色或深灰色、左对齐。
        # 逐 run/逐段校验，判据均为办公软件(PowerPoint/WPS)里真实生效的属性：
        #   - 字号取 run.font.size（磅）
        #   - 颜色取 run 实际字体颜色
        #   - 左对齐取段落 alignment=LEFT（未设置时办公软件默认左对齐）
        rec_shape = find_shape_by_text(slide, "相关记录：")
        ok = False
        msg = "缺失相关记录文本"
        if rec_shape is not None:
            paras = [p for p in rec_shape.text_frame.paragraphs if p.text.strip()]
            # 标题段（相关记录：）+ 项目符号段（沟通记录/证据材料/处理报告）都在校验范围内。
            title_ok = bool(paras) and norm_text(paras[0].text).startswith("相关记录")
            bullet_labels = ["沟通记录", "证据材料", "处理报告"]
            bullets_present = [w for w in bullet_labels if any(w in norm_text(p.text) for p in paras[1:])]
            content_ok = title_ok and len(bullets_present) == 3

            sizes = text_font_sizes(rec_shape)          # 标题+项目符号的全部文字字号
            colors = text_colors(rec_shape)             # 标题+项目符号的全部文字颜色
            size_ok = bool(sizes) and all(in_range(s, 15, 22, 0.5) for s in sizes)      # 字号15-22磅
            color_ok = bool(colors) and all(is_black_or_dark_gray(c) for c in colors)   # 黑色/深灰色
            align_ok = paragraph_align_left(rec_shape)                                  # 左对齐
            ok = content_ok and size_ok and color_ok and align_ok
            msg = f"标题+项目符号齐全={content_ok}，字号={sizes}({'OK' if size_ok else 'X'} 需15-22)，颜色={[rgb_hex(c) for c in colors]}({'OK' if color_ok else 'X'})，左对齐={align_ok}"
        self.add_point("+1 相关记录文本格式", 1, ok, msg)

    def check_line_styles(self):
        slide = self.slide
        # 细则：流程框、判断框和说明框——
        #  1) 轮廓均为“边线”(描边)；2) 线宽0.75-2磅；
        #  3) 蓝色框用深蓝色边线，绿色框用深绿色边线，橙色框用橙色边线。
        # 均为办公软件(PowerPoint/WPS)里真实生效的属性：<a:ln> 描边、线宽(EMU→磅)、线条颜色。
        labels = TOP_BLUE_LABELS + TOP_GREEN_LABELS + BOTTOM_LABELS
        details = []
        ok = True
        for label in labels:
            sh = find_shape_by_text(slide, label)
            if sh is None:
                ok = False
                details.append(f"{label}: 缺失")
                continue
            w = line_width_pt(sh)
            eff_w = effective_line_width_pt(sh)                     # 未显式设置时按默认 1.0pt 处理
            rgb = line_rgb(sh)
            edge_ok = has_edge_line(sh)                              # 轮廓为边线
            if label in TOP_BLUE_LABELS:
                color_ok = is_dark_blue(rgb)                        # 蓝色框→深蓝边线
            elif label in TOP_GREEN_LABELS:
                color_ok = is_teal_or_green(rgb)                    # 绿色框→深绿边线
            else:
                color_ok = is_orange(rgb)                           # 橙色框→橙色边线
            width_ok = in_range(eff_w, 0.75, 2.0)                   # 线宽0.75-2磅（未设置按默认1.0pt）
            this_ok = edge_ok and width_ok and color_ok
            ok = ok and this_ok
            details.append(f"{label}: 边线={edge_ok}, 线宽={'默认1.0' if w is None else round(w,2)}pt(需0.75-2), 颜色={rgb_hex(rgb)}, ok={this_ok}")

        # 说明框：左下、蓝色轮廓的矩形框（与记录说明框一致）。说明框轮廓为蓝色→深蓝边线。
        rec = None
        sw, shh = cm(self.prs.slide_width), cm(self.prs.slide_height)
        for sh in iter_shapes(slide):
            if is_rectangle(sh) and cm(sh.left) < sw * 0.35 and cm(sh.top) > shh * 0.55 and is_blue(line_rgb(sh)) and shape_area(sh) > 5:
                rec = sh
                break
        if rec is not None:
            w = line_width_pt(rec)
            eff_w = effective_line_width_pt(rec)                                            # 未显式设置按默认1.0pt
            rec_ok = has_edge_line(rec) and in_range(eff_w, 0.75, 2.0) and is_dark_blue(line_rgb(rec))
            ok = ok and rec_ok
            details.append(f"说明框: 边线={has_edge_line(rec)}, 线宽={'默认1.0' if w is None else round(w,2)}pt(需0.75-2), 颜色={rgb_hex(line_rgb(rec))}, ok={rec_ok}")
        else:
            ok = False
            details.append("说明框: 缺失")
        self.add_point("+3 流程框、判断框和说明框线条样式", 3, ok, "；".join(details[:8]) + ("；..." if len(details) > 8 else ""))

    def check_editability_score_point(self):
        slide = self.slide
        prs = self.prs
        page_area = cm(prs.slide_width) * cm(prs.slide_height)

        # 细则：用户可分别选中并修改/删减——
        #  1) 标题文本框；2) 每个流程框；3) 每个箭头；4) 每条回路线；
        #  5) 每个图标；6) 左下说明框中的文字。
        # 在办公软件(PowerPoint/WPS)里“可分别选中编辑”的充要条件是：这些内容都是各自独立的
        # 矢量对象（形状/文本框/连接线），而非被压平进一张位图。故逐类确认对象独立存在，
        # 并确认页面上没有承载这些内容的大位图（大位图=不可分别编辑）。

        # 1) 标题文本框：独立、含文本框。
        title = find_shape_by_text(slide, TARGET_TITLE)
        title_ok = title is not None and getattr(title, "has_text_frame", False)

        # 2) 每个流程框：全部流程框与判断框独立、且各自含可编辑文本框。
        box_labels = TOP_BLUE_LABELS + TOP_GREEN_LABELS + BOTTOM_LABELS
        box_list = [find_shape_by_text(slide, l) for l in box_labels]
        boxes_ok = all(b is not None and getattr(b, "has_text_frame", False) for b in box_list)

        # 3) 每个箭头 + 4) 每条回路线：均为连接线/线条对象（cxnSp/line），且不是图片。
        lines = get_lines(slide)
        # 主流程箭头(上6+下5=11) + 判断向下1 + 结案回流折线(多段) + 返工回流折线(多段)。
        # 这里以“线条对象足够多且都为矢量线条”确认箭头与回路线均可分别选中。
        lines_ok = len(lines) >= 11 and all(not is_picture(s) for s in lines)

        # 5) 每个图标：图标为独立可编辑矢量对象（形状/组合/含符号的文本），不是图片。
        #    统计页面上的小型非图片形状 + 含图标符号的文本对象。
        icon_syms = "🎧☎📞📄📋📃🗎📊📈👤🧑👨👩👥👫👬👭🤝📁🗂🗃📂🔍🔎⚙🛠🔧🛡✅✔☑☒☐✓❓？"
        icon_objs = 0
        for sh in iter_shapes(slide):
            if is_picture(sh) or is_line(sh):
                continue
            txt = text_of(sh)
            small = cm(sh.width) <= 1.6 and cm(sh.height) <= 1.6
            if small and not has_chinese(txt) and (is_oval(sh) or is_group(sh) or any(s in txt for s in icon_syms) or not txt.strip()):
                icon_objs += 1
            elif any(s in txt for s in icon_syms):
                icon_objs += 1
        icons_ok = icon_objs >= 1

        # 6) 左下说明框中的文字：独立文本框、含“相关记录：”文本。
        rec_text = find_shape_by_text(slide, "相关记录：")
        rec_ok = rec_text is not None and getattr(rec_text, "has_text_frame", False)

        # 反向约束：不存在承载核心内容的大位图（存在则说明被压平、无法分别编辑）。
        pictures = [s for s in iter_shapes(slide) if is_picture(s)]
        no_large_pic = not any(shape_area(p) > page_area * 0.10 for p in pictures)

        ok = title_ok and boxes_ok and lines_ok and icons_ok and rec_ok and no_large_pic
        self.add_point(
            "+3 承载页对象可分别选中编辑",
            3,
            ok,
            f"标题文本框={title_ok}；每个流程框={boxes_ok}；箭头/回路线对象数={len(lines)}(≥11:{lines_ok})；图标对象数={icon_objs}(可编辑:{icons_ok})；说明框文字={rec_ok}；无大位图={no_large_pic}",
        )

    def total_score(self) -> int:
        return sum(p.score for p in self.points if p.ok)

    def max_score(self) -> int:
        return sum(p.score for p in self.points)

    def print_report(self) -> None:
        # 保留兼容旧调试用途的打印接口：把与 build_result 一致的结果打印为可读文本。
        # 不用于对外主结果传递——主结果通过 evaluate() 的返回值 dict 提供。
        result = self.build_result()
        if result["status"] == "error":
            print(f"维度一：不通过（{result['error']}）")
            print("维度二：不检查")
            print("最终得分：0")
            return
        if not result["dim1_pass"]:
            print("维度一：不通过")
            print("维度二：不检查")
            print("最终得分：0")
            return
        hit_items = [it for it in result["dim2_items"] if it["hit"]]
        print("维度一：通过")
        print("维度二：评分细则结果")
        if not hit_items:
            print("无命中项")
        else:
            for it in hit_items:
                print(f"+{it['delta']}：{it['rule']}")
        print(f"最终得分：{result['total_score']} / {result['max_score']}")

    def build_result(
        self,
        script_id: str = "",
        file_name: str = "",
        status: str = "ok",
        error: Optional[str] = None,
    ) -> dict:
        """把内部评估状态汇总成统一约定的 dict 结构。"""
        if status == "error":
            return {
                "id": script_id,
                "file_name": file_name,
                "status": "error",
                "error": error,
                "dim1_pass": False,
                "dim1_reason": error or "",
                "dim2_items": [],
                "total_score": 0,
                "max_score": 0,
            }

        dim1_pass = bool(self.gate_results) and all(r.ok for r in self.gate_results)
        dim1_reason = "" if dim1_pass else "; ".join(f"{r.name}: {r.evidence}" for r in self.gate_results if not r.ok)

        dim2_items = []
        for p in self.points:
            rule_text = re.sub(r"^\s*[+-]?\d+\s*", "", p.name).lstrip("：: ")
            dim2_items.append({
                "rule": rule_text,
                "max_delta": p.score,
                "delta": p.score if p.ok else 0,
                "hit": bool(p.ok),
                "detail": "",
            })

        total_score = sum(it["delta"] for it in dim2_items)
        max_score = sum(it["max_delta"] for it in dim2_items)
        return {
            "id": script_id,
            "file_name": file_name,
            "status": "ok",
            "error": None,
            "dim1_pass": dim1_pass,
            "dim1_reason": dim1_reason,
            "dim2_items": dim2_items if dim1_pass else [],
            "total_score": total_score if dim1_pass else 0,
            "max_score": max_score,
        }


SCRIPT_ID = "056"
_SUPPORTED_EXTS = (".pptx",)


def _locate_target_file(dir_path: str) -> Optional[str]:
    """在给定目录内定位待评估 .pptx 文档；忽略以 ~$ 开头的临时文件。"""
    if not os.path.isdir(dir_path):
        return None
    candidates: list[str] = []
    for name in os.listdir(dir_path):
        if name.startswith("~$"):
            continue
        ext = os.path.splitext(name)[1].lower()
        if ext in _SUPPORTED_EXTS:
            candidates.append(name)
    if not candidates:
        return None
    candidates.sort()  # 稳定输出
    return os.path.join(dir_path, candidates[0])


def evaluate(dir_path: str) -> dict:
    """统一入口：接收“脚本所在目录的路径”，返回结构化评估结果。

    - 不做任何 stdout 副作用（不改 sys.stdout、不打印主结果）；
    - 不 sys.exit；异常统一收敛为 status="error"。
    """
    file_name = ""
    try:
        target = _locate_target_file(dir_path)
        if target is None:
            return Evaluator(dir_path).build_result(
                script_id=SCRIPT_ID,
                file_name="",
                status="error",
                error=f"未在目录中找到 .pptx 文件：{dir_path}",
            )
        file_name = os.path.basename(target)
        evaluator = Evaluator(target)
        if not evaluator.load():
            return evaluator.build_result(script_id=SCRIPT_ID, file_name=file_name)
        if evaluator.gate_check():
            evaluator.score_dimension2()
        return evaluator.build_result(script_id=SCRIPT_ID, file_name=file_name)
    except Exception as exc:  # 顶层兜底：任何未预期异常都转成 status=error。
        return Evaluator(dir_path).build_result(
            script_id=SCRIPT_ID,
            file_name=file_name,
            status="error",
            error=f"{type(exc).__name__}: {exc}",
        )


if __name__ == "__main__":
    # 仅用于本地调试：允许传入脚本所在目录的路径，默认使用脚本自身所在目录。
    _dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(_dir), ensure_ascii=False, indent=2))
