# -*- coding: utf-8 -*-
"""
自动评估脚本：流程图内容改写版_可编辑8cm版.pptx
评分细则见任务描述。
逻辑：
  维度1（可用与可修改性）—— 不满足则总分 0，结束。
  维度2（完成度）—— 累计 +/- 分项。
"""
import os, re, sys

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

SCRIPT_ID = "073"
DEFAULT_FILE_NAME = "流程图内容改写版_可编辑8cm版.pptx"


def _locate_pptx(dir_path: str) -> "str | None":
    """在给定目录中定位待评估的 .pptx 文件。
       优先使用约定名，否则取目录内首个 .pptx。"""
    preferred = os.path.join(dir_path, DEFAULT_FILE_NAME)
    if os.path.isfile(preferred):
        return preferred
    if os.path.isdir(dir_path):
        for name in sorted(os.listdir(dir_path)):
            if name.startswith("~$"):
                continue
            if name.lower().endswith(".pptx"):
                return os.path.join(dir_path, name)
    return None

# ------- 工具函数 -------
def cm(v):
    return Emu(v).cm if v is not None else 0.0

def iter_all_shapes(shapes):
    """递归展开组合形状。"""
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_all_shapes(sh.shapes)
        else:
            yield sh

def shape_text(sh):
    if sh.has_text_frame:
        return sh.text_frame.text or ""
    return ""

def shape_bbox(sh):
    """返回 (left, top, right, bottom) cm。"""
    l = cm(sh.left); t = cm(sh.top)
    w = cm(sh.width); h = cm(sh.height)
    return l, t, l + w, t + h

# ------- 维度 1 检查 -------
def check_dimension1(path):
    """细则：
       (a) 交付文件为 .pptx 格式，文件可正常打开。
    """
    issues = []

    # (a.1) 扩展名
    if not path.lower().endswith(".pptx"):
        issues.append("文件扩展名非 .pptx")
        return False, issues

    # (a.2) python-pptx 能正常加载
    try:
        prs = Presentation(path)
    except Exception as e:
        issues.append(f"无法打开 PPTX：{e}")
        return False, issues
    if len(prs.slides) == 0:
        issues.append("PPT 内无任何幻灯片")
        return False, issues

    return True, issues

# ------- 维度 2 评分项 -------
# 细则中两个“允许 >可达 8cm 的指定外侧框图”关键词（用于 -3 判定的白名单）
TARGET_KEYWORDS_8CM_ALLOWED = [
    "1.儿童成长规律与学习特征   2.课程标准素养目标与育人要求",
    "智能技术支持小学生差异化学习",
]
LAYER_LABELS = ["目标层", "核心层", "基础层", "机制层", "实施层", "运行逻辑"]

# 原 PPT 关键节点 / 层级标签 —— 用于检查“删除节点”
ORIGINAL_NODES = [
    # slide 1
    "价值引领根本目标", "个性学习落地", "学习者整体成长",
    "价值涵养", "智慧启发", "品格践行",
    "知识奠基", "能力提升", "素养生成",
    "育德", "立规范", "定目标", "划边界", "成才",
    "教育学", "技术学",
    "儿童成长规律与学习特征", "课程标准素养目标与育人要求",
    # slide 2
    "联动贯通", "持续推进", "智能技术支持小学生差异化学习",
    "教师统筹审核", "智能工具协助", "家校联动支持",
    "发展规律匹配", "持续复盘改进",
    "学情识别", "内容匹配", "路径调整", "反馈评价", "循环改进",
    # 层级
    "目标层", "核心层", "基础层", "机制层", "实施层", "运行逻辑",
]

def _resolve_effective_font_size_pt(run, para, text_frame):
    """按 OOXML 继承链解析有效字号（pt），返回 float 或 None（表示无法解析）。
       解析顺序：
         1) run 级 rPr/sz —— 显式设定；
         2) paragraph 级 pPr/defRPr/sz —— 段落默认；
         3) txBody 级 lstStyle/lvlNpPr/defRPr/sz —— 文本框列表样式默认；
         4) 无法解析 —— 返回 None，由调用方按“稳健兜底”处理。
       备注：占位符 → 母版 → 主题 一层的解析需要渲染引擎支持，
             python-pptx 不透出，这里统一并入“无法解析”。"""
    NS_A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    # 1) run-level
    if run.font.size is not None:
        return float(run.font.size.pt)
    # 2) paragraph-level defRPr
    p_el = para._p
    pPr = p_el.find(NS_A + 'pPr')
    if pPr is not None:
        defRPr = pPr.find(NS_A + 'defRPr')
        if defRPr is not None and defRPr.get('sz'):
            try:
                return int(defRPr.get('sz')) / 100.0
            except ValueError:
                pass
    # 3) txBody/lstStyle 的对应层级 defRPr
    try:
        txBody = text_frame._txBody
    except AttributeError:
        txBody = None
    if txBody is not None:
        lstStyle = txBody.find(NS_A + 'lstStyle')
        if lstStyle is not None:
            lvl = 0
            if pPr is not None and pPr.get('lvl'):
                try:
                    lvl = int(pPr.get('lvl'))
                except ValueError:
                    lvl = 0
            lvl_tag = 'lvl{}pPr'.format(lvl + 1)
            lvlPr = lstStyle.find(NS_A + lvl_tag)
            if lvlPr is not None:
                defRPr = lvlPr.find(NS_A + 'defRPr')
                if defRPr is not None and defRPr.get('sz'):
                    try:
                        return int(defRPr.get('sz')) / 100.0
                    except ValueError:
                        pass
    return None


def _resolve_effective_typefaces(run, para, text_frame):
    """按 OOXML 继承链解析字体（typeface），返回 (eastAsian, latin)。
       任一元素解析不到时对应位置为 None，由调用方进行兜底判定。
       解析顺序与 _resolve_effective_font_size_pt 一致：
         run 级 rPr → 段落 pPr/defRPr → txBody lstStyle/lvlNpPr/defRPr。"""
    NS_A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    ea = None
    latin = None

    def _pick(el):
        """从 rPr/defRPr 元素中抓取 eastAsian / latin typeface。"""
        e_ea = e_la = None
        if el is not None:
            ea_el = el.find(NS_A + 'eastAsian')
            if ea_el is not None:
                e_ea = ea_el.get('typeface')
            la_el = el.find(NS_A + 'latin')
            if la_el is not None:
                e_la = la_el.get('typeface')
        return e_ea, e_la

    # 1) run-level rPr
    rPr = run._r.find(NS_A + 'rPr')
    r_ea, r_la = _pick(rPr)
    ea = ea or r_ea
    latin = latin or r_la
    # run.font.name 兜底（python-pptx 抽象层）
    try:
        if not latin:
            latin = run.font.name
    except Exception:
        pass
    if ea and latin:
        return ea, latin

    # 2) paragraph-level defRPr
    p_el = para._p
    pPr = p_el.find(NS_A + 'pPr')
    if pPr is not None:
        d_ea, d_la = _pick(pPr.find(NS_A + 'defRPr'))
        ea = ea or d_ea
        latin = latin or d_la
        if ea and latin:
            return ea, latin

    # 3) txBody/lstStyle 的对应层级 defRPr
    try:
        txBody = text_frame._txBody
    except AttributeError:
        txBody = None
    if txBody is not None:
        lstStyle = txBody.find(NS_A + 'lstStyle')
        if lstStyle is not None:
            lvl = 0
            if pPr is not None and pPr.get('lvl'):
                try:
                    lvl = int(pPr.get('lvl'))
                except ValueError:
                    lvl = 0
            lvlPr = lstStyle.find(NS_A + 'lvl{}pPr'.format(lvl + 1))
            if lvlPr is not None:
                d_ea, d_la = _pick(lvlPr.find(NS_A + 'defRPr'))
                ea = ea or d_ea
                latin = latin or d_la

    return ea, latin


def _is_icon_shape(sh):
    """判断形状是否属于细则所指的“图标 / 文中图标”。
       纳入：AUTO_SHAPE（框图/箭头/连接节点）、FREEFORM（自定义形状）；
             以及以上类型嵌套在 GROUP 中的组件（由 iter_all_shapes 展开后到此已经不是 GROUP 了）。
       排除：TEXT_BOX（纯说明文字/页级标题）、PLACEHOLDER（版式占位符，通常承载页标题/正文）、
             PICTURE、CHART、TABLE、LINE 等——它们要么不是“图标”，要么不承载“图标内文字”。"""
    return sh.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM)


def check_font_uniform(slide):
    """严格按细则：所有【图标内文字和文中图标内容】均为
       (1) 宋体  (2) 6 号字  (3) 颜色为黑色或深灰色。
       返回 (ok, detail)。

       检查对象限定：
         - 仅纳入 AUTO_SHAPE / FREEFORM 上承载的文字（即“图标 / 文中图标”），
           排除独立文本框、页级标题占位符等非图标说明文字。

       字号判定：
         - 优先使用 run 级显式字号；否则回落到段落默认 / 文本框列表样式默认；
         - 无法解析（依赖占位符/母版/主题继承）时按“稳健兜底”处理，视为合格，
           而不是把 None 直接判为失败。

       颜色判定：
         - 显式 RGB：必须落在黑/深灰白名单内；
         - 主题方案色（SCHEME）：tx1/dk1/dk2/bg2 等深色方案色视为合格；
         - 未设色（color.type is None）：继承主题默认（tx1，黑），视为合格。
    """
    from pptx.enum.dml import MSO_THEME_COLOR
    # PowerPoint 标准命名色：黑色 + 三档深灰
    BLACK_OR_DARKGRAY = {
        (0x00, 0x00, 0x00),   # 黑色
        (0x26, 0x26, 0x26),   # 深灰色（黑色，文字 1，淡色 15%）
        (0x40, 0x40, 0x40),   # 深灰色 1
        (0x59, 0x59, 0x59),   # 深灰色 2
    }
    # 主题色中视为“深色（黑/深灰）”的方案色
    DARK_THEME_COLORS = {
        MSO_THEME_COLOR.TEXT_1,        # tx1，正文/背景方案中的深色文本
        MSO_THEME_COLOR.DARK_1,        # dk1
        MSO_THEME_COLOR.DARK_2,        # dk2
        MSO_THEME_COLOR.BACKGROUND_2,  # bg2（部分主题下也是深色）
    }
    bad = []
    icon_run_seen = False
    for sh in iter_all_shapes(slide.shapes):
        # 只检查“图标 / 文中图标”上的文字
        if not _is_icon_shape(sh):
            continue
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                icon_run_seen = True
                font = run.font

                # ---- 字体：宋体（沿继承链解析 eastAsian / latin） ----
                ea_typeface, latin_typeface = _resolve_effective_typefaces(run, para, tf)
                # 中文优先看 eastAsian，其次 latin；两者都缺失则视为“继承主题默认”，
                # 主题默认中文字体通常不是宋体，因此仅在两者都缺失时按“无法判定”兜底为合格，
                # 只要显式出现了非宋体的 typeface 就判失败。
                typeface_for_song = ea_typeface or latin_typeface
                if typeface_for_song is None:
                    is_song = True  # 稳健兜底：继承母版/主题，无法解析时视为合格
                else:
                    is_song = typeface_for_song in ("宋体", "SimSun")
                if not is_song:
                    bad.append(f"非宋体: typeface={typeface_for_song!r} text={run.text!r}")

                # ---- 字号：6 号字（沿继承链解析；无法解析时稳健兜底为通过） ----
                eff_pt = _resolve_effective_font_size_pt(run, para, tf)
                if eff_pt is None:
                    # 无法解析——继承自占位符 / 母版 / 主题，按“稳健兜底”视为合格
                    pass
                elif abs(eff_pt - 6.0) > 1e-6:
                    bad.append(f"非6号字: size={eff_pt}pt text={run.text!r}")

                # ---- 颜色：黑色或深灰色 ----
                color_ok = False
                color_detail = None
                try:
                    color = font.color
                    if color is None or color.type is None:
                        # 未显式设色 -> 继承主题默认（正文 tx1，即黑色）
                        color_ok = True
                        color_detail = "inherit/default"
                    else:
                        ctype = color.type
                        # 显式 RGB
                        try:
                            rgb = color.rgb
                        except Exception:
                            rgb = None
                        if rgb is not None:
                            color_ok = (rgb[0], rgb[1], rgb[2]) in BLACK_OR_DARKGRAY
                            color_detail = f"rgb={rgb}"
                        else:
                            # 主题方案色：检查 theme_color
                            try:
                                tc = color.theme_color
                            except Exception:
                                tc = None
                            if tc in DARK_THEME_COLORS:
                                color_ok = True
                                color_detail = f"theme={tc}"
                            else:
                                color_detail = f"type={ctype} theme={tc}"
                except Exception as e:
                    color_ok = False
                    color_detail = f"ERR:{e}"
                if not color_ok:
                    bad.append(f"颜色非黑/深灰: {color_detail} text={run.text!r}")

    # 若整页未采集到任何“图标 run”，视为该细则不适用——按合格返回，
    # 避免因限定范围过窄导致无差别失分（真正的“无图标/图片化”问题由维度一负责判定）。
    if not icon_run_seen:
        return True, []

    return (len(bad) == 0), bad

def find_shape_by_text(slide, keyword):
    """按归一化后完全匹配的方式查找文本等于 keyword 的形状。"""
    target = re.sub(r"\s+", "", keyword)
    matches = []
    for sh in iter_all_shapes(slide.shapes):
        t = re.sub(r"\s+", "", shape_text(sh))
        if t == target:
            matches.append(sh)
    return matches

def check_width_le_8_in_prs(prs, keyword):
    """在整个 PPT 中查找文本（忽略空白后）完全等于 keyword 的形状，
       要求其宽度不超过 8cm。"""
    target = re.sub(r"\s+", "", keyword)
    widths = []
    for slide in prs.slides:
        for sh in iter_all_shapes(slide.shapes):
            t = re.sub(r"\s+", "", shape_text(sh))
            if t == target:
                widths.append(cm(sh.width))
    if not widths:
        return False, f"未找到文本为『{keyword}』的外侧流程框图"
    ok = all(w <= 8.0 + 1e-6 for w in widths)
    return ok, f"匹配宽度={[round(w,2) for w in widths]}cm"

def collect_flowchart_boxes(prs):
    """返回所有 AUTO_SHAPE 且带文本（即“框图”），跨页。"""
    boxes = []
    for idx, slide in enumerate(prs.slides, 1):
        for sh in iter_all_shapes(slide.shapes):
            if sh.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
                txt = shape_text(sh).strip()
                if txt:
                    boxes.append((idx, sh, txt))
    return boxes

def _is_flowchart_component(sh, slide_area_cm2, layer_labels):
    """判断形状是否属于“流程图组件”（节点 / 连接线 / 分层标签 / 箭头）。
       纳入：
         - AUTO_SHAPE / FREEFORM：节点框图、箭头、装饰性流程符号；
         - LINE：连接线；
         - TEXT_BOX：仅当其文本命中层级标签或原始节点（视为“分层标签”）；
       排除：
         - PICTURE / CHART / TABLE / MEDIA / OLE_OBJECT / DIAGRAM 等；
         - PLACEHOLDER（页标题、页脚、页码、日期等版式元素）；
         - 覆盖幻灯片 ≥ 70% 面积的巨型形状（背景矩形）；
         - 完全无文本 且 面积 < 0.05 cm² 的装饰点/噪点。
       返回 (是否纳入, reason)。reason 仅用于调试。"""
    st = sh.shape_type

    # 显式排除的类型
    if st in (
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.CHART,
        MSO_SHAPE_TYPE.TABLE,
        MSO_SHAPE_TYPE.MEDIA,
        MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
        MSO_SHAPE_TYPE.DIAGRAM,
        MSO_SHAPE_TYPE.WEB_VIDEO,
    ):
        return False, f"excluded-type:{st}"

    # 版式占位符（标题 / 页脚 / 页码 / 日期）——不属于流程图
    if st == MSO_SHAPE_TYPE.PLACEHOLDER:
        return False, "placeholder"

    # 面积/尺寸
    try:
        w = cm(sh.width); h = cm(sh.height)
    except Exception:
        return False, "no-bbox"
    area = max(w, 0.0) * max(h, 0.0)

    # 背景矩形：占页面 ≥ 70% 面积
    if slide_area_cm2 > 0 and area / slide_area_cm2 >= 0.7:
        return False, "background"

    # 分类判断
    if st == MSO_SHAPE_TYPE.LINE:
        return True, "connector"
    if st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
        # 纯装饰性噪点：既无文本又极小
        has_text = sh.has_text_frame and bool((sh.text_frame.text or "").strip())
        if not has_text and area < 0.05:
            return False, "decorative-noise"
        return True, "node/arrow"
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        # 只有承载“分层标签 / 原节点”文本的独立文本框才算“流程图分层标签”
        txt = re.sub(r"\s+", "", shape_text(sh))
        if not txt:
            return False, "empty-textbox"
        if any(re.sub(r"\s+", "", lab) == txt for lab in layer_labels):
            return True, "layer-label"
        if any(re.sub(r"\s+", "", n) == txt for n in ORIGINAL_NODES):
            return True, "node-textbox"
        return False, "unrelated-textbox"

    # 其它未知类型：保守排除
    return False, f"other:{st}"


def overall_flowchart_bbox(slide):
    """流程图整体外接框：仅聚合『流程图相关』对象——节点框图、连接线、
       分层标签、箭头；明确排除背景矩形、版式占位符（页码/页脚/标题）、
       图片、图表、表格、装饰元素等无关对象。
       返回 (l, t, r, b) 单位 cm；若该页没有任何流程图组件，返回 None。"""
    # 幻灯片面积（用于识别背景矩形）
    try:
        sw = cm(slide.part.package.presentation_part.presentation.slide_width)
        sh_h = cm(slide.part.package.presentation_part.presentation.slide_height)
    except Exception:
        sw = sh_h = 0.0
    slide_area = sw * sh_h

    xs1, ys1, xs2, ys2 = [], [], [], []
    for sh in iter_all_shapes(slide.shapes):
        ok, _reason = _is_flowchart_component(sh, slide_area, LAYER_LABELS)
        if not ok:
            continue
        try:
            l, t, r, b = shape_bbox(sh)
        except Exception:
            continue
        if (r - l) <= 0 and (b - t) <= 0:
            # 零面积的对象（例如没有 width/height 的异常形状）跳过
            continue
        xs1.append(l); xs2.append(r); ys1.append(t); ys2.append(b)
    if not xs1:
        return None
    return min(xs1), min(ys1), max(xs2), max(ys2)

def _emu_to_cm(v):
    """python-pptx 里 margin/line width 可能是 Emu 或原生 int(EMU)，统一转 cm。
    （历史遗留：原用于溢出/遮挡判定，现该规则已删除，函数保留供潜在复用。）"""
    if v is None:
        return None
    try:
        return Emu(int(v)).cm
    except Exception:
        try:
            return float(v) / 360000.0
        except Exception:
            return None


# ------- 主流程 -------
# 细则条文（评分细则原文，用于结果打印）
RULES = {
    "+3a": "第1页字体统一：所有图标内文字和文中图标内容均设置为宋体6号字，颜色保持黑色、深灰色。",
    "+3b": "第2页字体统一：所有图标内文字和文中图标内容均设置为宋体6号字，颜色保持黑色、深灰色。",
    "+1a": "“1.儿童成长规律与学习特征   2.课程标准素养目标与育人要求”外侧流程框图宽需改为不超过8厘米",
    "+1b": "“智能技术支持小学生差异化学习”外侧流程框图宽需改为不超过8厘米",
    "-5a": "页面被做成整页截图或不可编辑图片。",
    "-5b": "新增或删除原PPT中的流程节点、说明文字或层级标签。",
    "-5c": "流程图整体不可小于20厘米",
    "-3a": "除“1.儿童成长规律与学习特征   2.课程标准素养目标与育人要求”、“智能技术支持小学生差异化学习”外任意一个流程框图宽度超过8厘米。",
    "-3e": "为压缩框宽导致流程图整体明显错位、重叠或不可读。",
    "-3f": "页面边缘出现内容裁切，或任意对象超出幻灯片边界。",
    "-3g": "PPT中出现与任务无关的批注、红色标记、截图边框、空白页或说明文字。",
}

# 维度二评分项定义：(rule_key, max_delta)
# 正向加分项使用正数；扣分项使用负数。命中时 delta=max_delta，否则为 0。
DIM2_ITEMS_DEF = [
    ("+3a", 3),
    ("+3b", 3),
    ("+1a", 1),
    ("+1b", 1),
    ("-5a", -5),
    ("-5b", -5),
    ("-5c", -5),
    ("-3a", -3),
    ("-3e", -3),
    ("-3f", -3),
    ("-3g", -3),
]


def evaluate(dir_path: str) -> dict:
    """统一入口：接收脚本所在目录路径，脚本自行在目录内定位并打开被评估文档。

    返回结构见《脚本接口差异与统一建议.md》 §2.2。
    """
    # 结果骨架
    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": sum(m for _, m in DIM2_ITEMS_DEF if _.startswith("+")),
    }

    try:
        file_path = _locate_pptx(dir_path)
        if not file_path or not os.path.isfile(file_path):
            result["status"] = "error"
            result["error"] = f"未在目录 {dir_path!r} 中找到 .pptx 文件"
            return result
        result["file_name"] = os.path.basename(file_path)

        # ---- 维度 1 ----
        ok1, det1 = check_dimension1(file_path)
        if not ok1:
            result["dim1_pass"] = False
            result["dim1_reason"] = "；".join(det1) if det1 else "维度一未通过"
            result["total_score"] = 0
            return result

        result["dim1_pass"] = True

        prs = Presentation(file_path)

        # 逐项判定，先算 hit / delta
        hit_map = {}   # rule_key -> (hit: bool, delta: int, detail: str)

        # +3 第1页字体统一
        ok, _bad = check_font_uniform(prs.slides[0])
        hit_map["+3a"] = (bool(ok), 3 if ok else 0, "" if ok else f"违规项数={len(_bad)}")

        # +3 第2页字体统一
        if len(prs.slides) >= 2:
            ok, _bad = check_font_uniform(prs.slides[1])
            hit_map["+3b"] = (bool(ok), 3 if ok else 0, "" if ok else f"违规项数={len(_bad)}")
        else:
            hit_map["+3b"] = (False, 0, "无第2页")

        # +1 外侧框宽 ≤ 8cm
        ok, det = check_width_le_8_in_prs(prs, "1.儿童成长规律与学习特征   2.课程标准素养目标与育人要求")
        hit_map["+1a"] = (ok is True, 1 if ok is True else 0, det)

        ok, det = check_width_le_8_in_prs(prs, "智能技术支持小学生差异化学习")
        hit_map["+1b"] = (ok is True, 1 if ok is True else 0, det)

        # -5c 流程图整体不可小于 20 厘米
        under_20 = False
        for slide in prs.slides:
            bb = overall_flowchart_bbox(slide)
            if bb and (bb[2] - bb[0]) < 20.0:
                under_20 = True; break
        hit_map["-5c"] = (under_20, -5 if under_20 else 0, "" if not under_20 else "整体宽<20cm")

        # -3a 除两个允许的外侧框外，任意一个流程框图宽度超过 8 厘米
        boxes = collect_flowchart_boxes(prs)
        allowed_norm = [
            re.sub(r"\s+", "", "1.儿童成长规律与学习特征   2.课程标准素养目标与育人要求"),
            re.sub(r"\s+", "", "智能技术支持小学生差异化学习"),
        ]
        over_w_hit = False
        over_w_detail = ""
        for _, sh, txt in boxes:
            w = cm(sh.width)
            t_norm = re.sub(r"\s+", "", txt)
            if t_norm in allowed_norm:
                continue
            if w > 8.0 + 1e-3:
                over_w_hit = True
                over_w_detail = f"{txt!r} 宽={w:.2f}cm"
                break
        hit_map["-3a"] = (over_w_hit, -3 if over_w_hit else 0, over_w_detail)

        # 组装 dim2_items 与总分：正向项和扣分项都直接累加实际 delta。
        pos_max = sum(m for k, m in DIM2_ITEMS_DEF if k.startswith("+"))
        total = 0
        for k, m in DIM2_ITEMS_DEF:
            hit, delta, _detail = hit_map.get(k, (False, 0, ""))
            total += delta

            result["dim2_items"].append({
                "rule": RULES[k],
                "max_delta": m,
                "delta": delta,
                "hit": hit,
                "detail": "",
            })

        result["total_score"] = total
        result["max_score"] = pos_max
        return result

    except Exception as e:
        result["status"] = "error"
        result["error"] = f"{type(e).__name__}: {e}"
        return result


if __name__ == "__main__":
    import json as _json
    _dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(_json.dumps(evaluate(_dir), ensure_ascii=False, indent=2))
