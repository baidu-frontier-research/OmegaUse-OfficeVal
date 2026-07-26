#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""自动评估“豆包制作PPT001_修改版(5).pptx”。

评分逻辑：
1. 先检查维度1：交付文件为.pptx格式，文件可正常打开。
   不满足时直接 0 分，不再检查维度2。
2. 满足维度1后，逐条自动检测维度2的得分点和扣分点；命中则累计对应分值。
3. 通过 evaluate(dir_path) 返回结构化字典，供批量运行器汇总。

依赖：python-pptx、Pillow、lxml（通常安装 python-pptx 时会带上 lxml）。
"""

from __future__ import annotations

import json
import math
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Optional

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL

EMU_PER_PT = 12700

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


@dataclass
class CheckResult:
    score: int
    name: str
    passed: bool
    detail: str


@dataclass
class Rect:
    left: int
    top: int
    width: int
    height: int

    @property
    def right(self) -> int:
        return self.left + self.width

    @property
    def bottom(self) -> int:
        return self.top + self.height

    @property
    def cx(self) -> float:
        return self.left + self.width / 2

    @property
    def cy(self) -> float:
        return self.top + self.height / 2

    @property
    def area(self) -> int:
        return max(0, self.width) * max(0, self.height)


def shape_text(shape) -> str:
    if not hasattr(shape, "text"):
        return ""
    return (shape.text or "").strip()


def normalized_text(text: str) -> str:
    return "".join((text or "").split())


def all_slide_text(slide) -> str:
    return "\n".join(shape_text(shape) for shape in slide.shapes if shape_text(shape))


def iter_picture_shapes(slide):
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape


def rect_of(shape) -> Rect:
    return Rect(int(shape.left), int(shape.top), int(shape.width), int(shape.height))


def overlap_area(a: Rect, b: Rect) -> int:
    x1 = max(a.left, b.left)
    y1 = max(a.top, b.top)
    x2 = min(a.right, b.right)
    y2 = min(a.bottom, b.bottom)
    return max(0, x2 - x1) * max(0, y2 - y1)


def rgb_tuple_from_hex(value: str) -> tuple[int, int, int]:
    value = value.strip().lstrip("#")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def color_distance(c1: tuple[int, int, int], c2: tuple[int, int, int]) -> float:
    return math.sqrt(sum((a - b) ** 2 for a, b in zip(c1, c2)))


def classify_green(hex_color: str | None) -> Optional[str]:
    """把颜色粗略归类为 deep_green / light_green / other。

    细则只给出“深绿色/浅绿色”的语义，没有限定精确色值。这里采用自动化近似：
    - 绿色通道必须明显占优；
    - 亮度低于阈值视为深绿色，较高视为浅绿色。
    """
    if not hex_color:
        return None
    r, g, b = rgb_tuple_from_hex(hex_color)
    if g < 70 or g < max(r, b) + 10:
        return "other"
    brightness = 0.299 * r + 0.587 * g + 0.114 * b
    return "deep_green" if brightness < 150 else "light_green"


def is_blue(hex_color: str | None) -> bool:
    """判断颜色是否为蓝色系。

    细则只要求“蓝色填充”，未限定精确色值，故按颜色语义近似判断：
    蓝色通道明显占优即视为蓝色。
    """
    if not hex_color:
        return False
    r, g, b = rgb_tuple_from_hex(hex_color)
    return b >= 100 and b > r + 20 and b > g + 20


def is_white(hex_color: str | None, tolerance: float = 40.0) -> bool:
    """判断颜色是否为白色（含近白色容差）。细则要求“白色外边”。"""
    if not hex_color:
        return False
    return color_distance(rgb_tuple_from_hex(hex_color), (255, 255, 255)) <= tolerance


def shape_fill_hex(shape) -> Optional[str]:
    if not hasattr(shape, "fill"):
        return None
    try:
        fill = shape.fill
        if fill.type != MSO_FILL.SOLID:
            return None
        rgb = fill.fore_color.rgb
        return str(rgb).upper() if rgb else None
    except Exception:
        return None


def get_slide_xml_root(pptx_path: Path, slide_number: int):
    with zipfile.ZipFile(pptx_path) as zf:
        xml = zf.read(f"ppt/slides/slide{slide_number}.xml")
    return etree.fromstring(xml)


def shape_id(shape) -> Optional[str]:
    try:
        return str(shape._element.xpath(".//p:cNvPr", namespaces=NS)[0].get("id"))
    except Exception:
        return None


def find_xml_shape_by_id(root, sid: str):
    matches = root.xpath(f".//p:sp[p:nvSpPr/p:cNvPr[@id='{sid}']]", namespaces=NS)
    return matches[0] if matches else None


def text_shape_xml_nodes_with_text(root, target: str):
    nodes = []
    for sp in root.xpath(".//p:sp", namespaces=NS):
        text = "".join(sp.xpath(".//a:t/text()", namespaces=NS))
        if target in text:
            nodes.append(sp)
    return nodes


def _load_theme_colors(pptx_path: Path) -> dict[str, str]:
    """加载 theme1.xml 的方案色映射：`accent1`/`dk1`/`lt1`/... → 六位 hex。

    PPTX 中文字/线条颜色可能写成 ``<a:schemeClr val="accent1"/>`` 或
    ``<a:sysClr lastClr="FFFFFF" val="window"/>``，只解析 ``srgbClr`` 会漏判。
    """
    try:
        with zipfile.ZipFile(str(pptx_path)) as zf:
            xml = zf.read("ppt/theme/theme1.xml")
        root = etree.fromstring(xml)
    except Exception:
        return {}
    colors: dict[str, str] = {}
    for elem in root.xpath(".//a:clrScheme/*", namespaces=NS):
        name = etree.QName(elem.tag).localname
        for child in elem:
            tag = etree.QName(child.tag).localname
            if tag == "srgbClr":
                val = child.get("val")
                if val:
                    colors[name] = val.upper()
                    break
            if tag == "sysClr":
                val = child.get("lastClr")
                if val:
                    colors[name] = val.upper()
                    break
    return colors


# schemeClr 中 dk1/lt1/dk2/lt2 与 tx1/bg1/tx2/bg2 互为别名（Office Open XML §20.1.4.1.24）。
_SCHEME_ALIASES = {"bg1": "lt1", "tx1": "dk1", "bg2": "lt2", "tx2": "dk2",
                   "lt1": "bg1", "dk1": "tx1", "lt2": "bg2", "dk2": "tx2"}


def _resolve_color_hex(clr_container, theme_colors: dict[str, str]) -> Optional[str]:
    """从 ``<a:solidFill>``/``<a:ln>`` 之类的容器里解析出 hex 颜色。

    识别 ``srgbClr``（直接色）、``schemeClr``（主题方案色）、``sysClr``（系统色 lastClr）。
    未处理 lumMod/lumOff 等修饰——白/蓝的语义判定对轻微亮度调整仍能容忍。
    """
    if clr_container is None:
        return None
    for child in clr_container:
        tag = etree.QName(child.tag).localname
        if tag == "srgbClr":
            val = child.get("val")
            return val.upper() if val else None
        if tag == "schemeClr":
            name = child.get("val") or ""
            hex_val = theme_colors.get(name)
            if hex_val is None and name in _SCHEME_ALIASES:
                hex_val = theme_colors.get(_SCHEME_ALIASES[name])
            return hex_val
        if tag == "sysClr":
            val = child.get("lastClr")
            return val.upper() if val else None
    return None


def _run_has_wordart_features(rpr, sp) -> tuple[bool, list[str]]:
    """判断某个 run（``a:rPr``）及其所属形状是否具备 WordArt 特征。

    仅"填充为蓝色 + 轮廓为白色"并不足以证明是艺术字——普通文本框也可以通过
    "文本填充/文本轮廓"设置同样的颜色。因此这里额外核查真正的艺术字标志位：

    - ``a:rPr/a:ln``：文本轮廓本身即艺术字属性（普通文字段落没有 ``a:ln`` 子节点）。
    - ``a:rPr/a:effectLst``：文字效果（发光/阴影/映像/柔化边缘）。
    - ``a:rPr/a:uFill``/``a:uLn``/``a:highlight``：进阶文字效果。
    - ``a:rPr/a:scene3d``、``a:sp3d``：3D 文字效果。
    - 形状 ``p:txBody/a:bodyPr/a:prstTxWarp``（非 ``textNoShape``）：文字变形（艺术字最强指征）。
    - 形状 ``p:style/a:fontRef``：主题级 WordArt 样式引用。
    """
    reasons: list[str] = []
    if rpr.xpath("./a:ln", namespaces=NS):
        reasons.append("rPr/ln 文本轮廓")
    if rpr.xpath("./a:effectLst/*", namespaces=NS):
        reasons.append("rPr/effectLst 文字效果")
    if rpr.xpath("./a:uFill | ./a:uLn | ./a:highlight", namespaces=NS):
        reasons.append("rPr 高级文字装饰")
    if rpr.xpath("./a:scene3d | ./a:sp3d", namespaces=NS):
        reasons.append("rPr 3D 文字效果")

    warps = sp.xpath(".//p:txBody/a:bodyPr/a:prstTxWarp/@prst", namespaces=NS)
    if any(w and w != "textNoShape" for w in warps):
        reasons.append("bodyPr/prstTxWarp 文字变形")
    if sp.xpath(".//p:style/a:fontRef", namespaces=NS):
        reasons.append("style/fontRef 主题艺术字样式")

    return bool(reasons), reasons


def has_wordart_white_line_blue_fill(sp, theme_colors: dict[str, str]) -> tuple[bool, str]:
    """检查文字 run 是否为"白色外边（轮廓线）+ 蓝色填充"的艺术字。

    组合判定：同一 run 既要满足颜色语义（蓝色文本填充 + 白色文本轮廓），
    又要具备至少一项 WordArt 结构特征（详见 ``_run_has_wordart_features``）；
    避免"普通文本框设置同样颜色"误判为艺术字。颜色解析支持 srgbClr / schemeClr / sysClr。
    """
    for rpr in sp.xpath(".//a:rPr", namespaces=NS):
        fill_hex = None
        for fill in rpr.xpath("./a:solidFill", namespaces=NS):
            fill_hex = _resolve_color_hex(fill, theme_colors)
            if fill_hex:
                break
        line_hex = None
        ln_nodes = rpr.xpath("./a:ln", namespaces=NS)
        for ln in ln_nodes:
            for sf in ln.xpath("./a:solidFill", namespaces=NS):
                line_hex = _resolve_color_hex(sf, theme_colors)
                if line_hex:
                    break
            if line_hex:
                break

        fill_ok = fill_hex is not None and is_blue(fill_hex)
        line_ok = line_hex is not None and is_white(line_hex)
        if not (fill_ok and line_ok):
            continue

        is_wordart, reasons = _run_has_wordart_features(rpr, sp)
        if is_wordart:
            return True, (
                f"填充={fill_hex}(蓝色)，轮廓={line_hex}(白色)，"
                f"艺术字特征：{'、'.join(reasons)}"
            )
    return False, ""


def title_font_sizes_pt(root, target: str) -> list[float]:
    sizes = []
    for sp in text_shape_xml_nodes_with_text(root, target):
        for sz in sp.xpath(".//a:rPr/@sz", namespaces=NS):
            try:
                sizes.append(int(sz) / 100)
            except ValueError:
                pass
    return sizes


def _first_sz(elem, xpath: str) -> Optional[float]:
    """从 XML 子树按 xpath 找第一个 @sz（单位百分点），失败返回 None。"""
    if elem is None:
        return None
    try:
        vals = elem.xpath(xpath, namespaces=NS)
    except Exception:
        return None
    for v in vals:
        try:
            return int(v) / 100
        except (TypeError, ValueError):
            continue
    return None


def _run_explicit_size(run_elem) -> Optional[float]:
    return _first_sz(run_elem, "./a:rPr/@sz")


def _paragraph_default_size(paragraph_elem, level_zero_based: int) -> Optional[float]:
    """段落级默认字号：优先 a:pPr/a:defRPr/@sz，再向上到形状的 a:lstStyle/a:lvlXpPr。"""
    return _first_sz(paragraph_elem, "./a:pPr/a:defRPr/@sz")


def _lst_style_size(txBody_elem, level_zero_based: int) -> Optional[float]:
    """形状内 a:lstStyle 的默认字号（按段落级 0..8 对应 lvl1pPr..lvl9pPr）。"""
    if txBody_elem is None:
        return None
    lvl = max(0, min(8, level_zero_based)) + 1
    lvl_tag = f"a:lvl{lvl}pPr"
    if lvl == 1:
        candidates = [f"./a:lstStyle/{lvl_tag}/a:defRPr/@sz",
                      "./a:lstStyle/a:defPPr/a:defRPr/@sz"]
    else:
        candidates = [f"./a:lstStyle/{lvl_tag}/a:defRPr/@sz"]
    for xp in candidates:
        v = _first_sz(txBody_elem, xp)
        if v is not None:
            return v
    return None


def _placeholder_key(sp_elem) -> Optional[tuple[Optional[str], Optional[str]]]:
    """占位符匹配键：(type, idx)。二者都可能缺省——按 OOXML 规则视作 body/0。"""
    phs = sp_elem.xpath(".//p:nvSpPr/p:nvPr/p:ph", namespaces=NS)
    if not phs:
        return None
    ph = phs[0]
    return ph.get("type"), ph.get("idx")


def _find_matching_placeholder(container_root, key: tuple[Optional[str], Optional[str]]):
    """在 slideLayout/slideMaster 的根节点中匹配同 type/idx 的占位符 sp。"""
    if container_root is None:
        return None
    ph_type, ph_idx = key
    for sp in container_root.xpath(".//p:sp", namespaces=NS):
        phs = sp.xpath(".//p:nvSpPr/p:nvPr/p:ph", namespaces=NS)
        if not phs:
            continue
        t = phs[0].get("type")
        i = phs[0].get("idx")
        if t == ph_type and i == ph_idx:
            return sp
        # 常见互认：title / ctrTitle 都视为标题占位符；idx 缺省视为 "0"。
        if ph_type in ("title", "ctrTitle") and t in ("title", "ctrTitle") and (i or "0") == (ph_idx or "0"):
            return sp
    return None


def _master_style_size(master_root, ph_type: Optional[str], level_zero_based: int) -> Optional[float]:
    """slideMaster 的 titleStyle / bodyStyle / otherStyle 里查对应层级的默认字号。"""
    if master_root is None:
        return None
    if ph_type in ("title", "ctrTitle"):
        style_tag = "p:titleStyle"
    elif ph_type in ("body", "subTitle", None):
        style_tag = "p:bodyStyle"
    else:
        style_tag = "p:otherStyle"
    lvl = max(0, min(8, level_zero_based)) + 1
    return _first_sz(master_root, f".//{style_tag}/a:lvl{lvl}pPr/a:defRPr/@sz")


def _read_related_xml(pptx_path: Path, part_path: str):
    try:
        with zipfile.ZipFile(str(pptx_path)) as zf:
            return etree.fromstring(zf.read(part_path))
    except Exception:
        return None


def _resolve_slide_layout_and_master(pptx_path: Path, slide_number: int):
    """通过 rels 追踪 slideN → slideLayout → slideMaster，返回二者的 XML 根节点。"""
    rels_path = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
    try:
        with zipfile.ZipFile(str(pptx_path)) as zf:
            rels_root = etree.fromstring(zf.read(rels_path))
            layout_target = None
            for rel in rels_root.xpath(
                "./*[local-name()='Relationship' and contains(@Type,'slideLayout')]",
            ):
                layout_target = rel.get("Target")
                break
            if not layout_target:
                return None, None
            layout_part = "ppt/" + layout_target.replace("../", "")
            layout_root = etree.fromstring(zf.read(layout_part))
            layout_rels = layout_part.rsplit("/", 1)
            layout_rels_path = f"{layout_rels[0]}/_rels/{layout_rels[1]}.rels"
            master_root = None
            try:
                lrels = etree.fromstring(zf.read(layout_rels_path))
                for rel in lrels.xpath(
                    "./*[local-name()='Relationship' and contains(@Type,'slideMaster')]",
                ):
                    master_target = rel.get("Target")
                    if master_target:
                        master_part = "ppt/" + master_target.replace("../", "").lstrip("/")
                        # slideLayout 位于 ppt/slideLayouts/，master target 通常是 "../slideMasters/slideMasterN.xml"
                        if master_target.startswith("../"):
                            master_part = "ppt/" + master_target[3:]
                        elif master_target.startswith("/"):
                            master_part = master_target.lstrip("/")
                        master_root = etree.fromstring(zf.read(master_part))
                        break
            except KeyError:
                pass
            return layout_root, master_root
    except Exception:
        return None, None


def resolve_run_effective_size(
    pptx_path: Path, slide_number: int, sp_elem, paragraph_elem, run_elem,
) -> tuple[Optional[float], str]:
    """按 OOXML 继承链解析 run 的最终字号（磅）。

    继承顺序：run.rPr → paragraph.pPr.defRPr → shape.lstStyle.lvlXpPr →
    slideLayout 同占位符 → slideMaster 同占位符 → slideMaster titleStyle/bodyStyle/otherStyle。
    未在任一层找到即返回 (None, 说明)。返回值第二项用于诊断输出。
    """
    # 段落级号
    level = 0
    lvl_attr = paragraph_elem.xpath("./a:pPr/@lvl", namespaces=NS)
    if lvl_attr:
        try:
            level = int(lvl_attr[0])
        except (TypeError, ValueError):
            level = 0

    v = _run_explicit_size(run_elem)
    if v is not None:
        return v, "run.rPr@sz"

    v = _paragraph_default_size(paragraph_elem, level)
    if v is not None:
        return v, "pPr/defRPr@sz"

    txBody = sp_elem.xpath(".//p:txBody", namespaces=NS)
    if txBody:
        v = _lst_style_size(txBody[0], level)
        if v is not None:
            return v, f"shape.lstStyle.lvl{level + 1}pPr@sz"

    ph_key = _placeholder_key(sp_elem)
    layout_root, master_root = _resolve_slide_layout_and_master(pptx_path, slide_number)

    if ph_key is not None:
        # slideLayout 同占位符（rPr / pPr.defRPr / lstStyle）
        layout_sp = _find_matching_placeholder(layout_root, ph_key)
        if layout_sp is not None:
            v = _first_sz(layout_sp, f".//a:p[{level + 1}]/a:pPr/a:defRPr/@sz")
            if v is None:
                v = _first_sz(layout_sp, ".//a:p/a:pPr/a:defRPr/@sz")
            if v is None:
                lb = layout_sp.xpath(".//p:txBody", namespaces=NS)
                if lb:
                    v = _lst_style_size(lb[0], level)
            if v is not None:
                return v, "slideLayout(占位符)"
        # slideMaster 同占位符
        master_sp = _find_matching_placeholder(master_root, ph_key)
        if master_sp is not None:
            v = _first_sz(master_sp, f".//a:p[{level + 1}]/a:pPr/a:defRPr/@sz")
            if v is None:
                v = _first_sz(master_sp, ".//a:p/a:pPr/a:defRPr/@sz")
            if v is None:
                mb = master_sp.xpath(".//p:txBody", namespaces=NS)
                if mb:
                    v = _lst_style_size(mb[0], level)
            if v is not None:
                return v, "slideMaster(占位符)"
        # slideMaster titleStyle / bodyStyle / otherStyle
        v = _master_style_size(master_root, ph_key[0], level)
        if v is not None:
            return v, "slideMaster.titleStyle/bodyStyle"

    return None, "未在 slide/layout/master 中找到有效 @sz"


def collect_title_effective_sizes(
    pptx_path: Path, slide_number: int, target: str,
) -> list[tuple[float, str]]:
    """收集所有包含 target 文本的 run 的最终字号（含出处诊断）。

    若某个 run 未查到字号，用 (None, 原因) 也会记录，供上层判定"字号不明"。
    """
    root = get_slide_xml_root(pptx_path, slide_number)
    result: list[tuple[float, str]] = []
    unresolved: list[str] = []
    for sp in text_shape_xml_nodes_with_text(root, target):
        for paragraph in sp.xpath(".//a:p", namespaces=NS):
            para_text = "".join(paragraph.xpath(".//a:t/text()", namespaces=NS))
            # 只在段落级也含目标文本时才认为该段落隶属于标题；防止把同 shape 内其他段落
            # （如副标题、注释）的字号也拿来判定。
            if target not in para_text:
                continue
            for run in paragraph.xpath("./a:r", namespaces=NS):
                run_text = "".join(run.xpath(".//a:t/text()", namespaces=NS))
                if not run_text.strip():
                    continue
                size, origin = resolve_run_effective_size(pptx_path, slide_number, sp, paragraph, run)
                if size is not None:
                    result.append((size, f"[{run_text[:6]}…] {origin}={size}pt"))
                else:
                    unresolved.append(f"[{run_text[:6]}…] {origin}")
    # 把 unresolved 塞进返回结构：约定 size=-1 表示"未解析出字号"。
    for reason in unresolved:
        result.append((-1.0, reason))
    return result


def picture_aspect_ratio(shape) -> float:
    # 优先检查PPT中实际摆放尺寸；这比原图像素更符合“PPT图片横纵比”的评估意图。
    if shape.height:
        return float(shape.width) / float(shape.height)
    return 0.0


def is_square_picture(shape, tolerance: float = 0.03) -> bool:
    return abs(picture_aspect_ratio(shape) - 1.0) <= tolerance


def two_sides_balanced(left_rect: Rect, right_rect: Rect, tolerance: float = 0.08) -> bool:
    if min(left_rect.width, right_rect.width) <= 0:
        return False
    ratio = left_rect.width / right_rect.width
    return abs(ratio - 1.0) <= tolerance


def find_slide2_main_title_rect(slide) -> Optional[Rect]:
    """找到第二页顶部总标题“田间防护：护鸟带与驱鸟器应用说明”的矩形。"""
    target = normalized_text("田间防护：护鸟带与驱鸟器应用说明")
    for shape in slide.shapes:
        if target in normalized_text(shape_text(shape)):
            return rect_of(shape)
    return None


def find_panel_title_shapes(slide):
    bird_band = []
    repeller = []
    for shape in slide.shapes:
        text = normalized_text(shape_text(shape))
        # 顶部总标题同时包含“护鸟带”和“驱鸟器”，不能作为单个板块标题。
        if "护鸟带" in text and "驱鸟器" not in text:
            bird_band.append(shape)
        if "驱鸟器" in text and "护鸟带" not in text:
            repeller.append(shape)
    return bird_band, repeller


def infer_content_panels(slide, slide_w: int, slide_h: int) -> Optional[tuple[Rect, Rect, str]]:
    """推断第二页护鸟带/驱鸟器两个实际板块区域。

    只有大背景块"真正包含"对应标题时才会归属该板块——去除按中心距离的兜底，
    避免把同一个大背景块同时最近于两个标题、或把无关背景块误归属为板块。
    如果大背景块方案不成立，则退回按板块标题/正文文字包围盒推断。
    """
    bird_shapes, repeller_shapes = find_panel_title_shapes(slide)
    if not bird_shapes or not repeller_shapes:
        return None

    title_bird = rect_of(bird_shapes[0])
    title_repeller = rect_of(repeller_shapes[0])

    candidates = []
    for shape in slide.shapes:
        color = shape_fill_hex(shape)
        sr = rect_of(shape)
        if not color or sr.area == 0:
            continue
        if sr.width >= slide_w * 0.25 and sr.height >= slide_h * 0.18:
            candidates.append(sr)

    def containing_panel(title_rect: Rect) -> Optional[Rect]:
        containing = [r for r in candidates if overlap_area(r, title_rect) / max(1, title_rect.area) >= 0.50]
        if containing:
            return max(containing, key=lambda r: r.area)
        return None

    bird_panel = containing_panel(title_bird)
    repeller_panel = containing_panel(title_repeller)
    # 大背景块可能同时把两个标题都罩住（同一整块背景），此时视为无有效背景。
    if bird_panel and repeller_panel and bird_panel == repeller_panel:
        bird_panel = repeller_panel = None

    if not bird_panel or not repeller_panel:
        # 兜底：使用各自标题/正文文字包围盒（严格按各自的文字位置，不做最近邻猜测）。
        bird_panel = bounding_rect([rect_of(s) for s in bird_shapes])
        repeller_panel = bounding_rect([rect_of(s) for s in repeller_shapes])
        return bird_panel, repeller_panel, "未找到独立大背景块，按板块文字包围盒推断"

    return bird_panel, repeller_panel, "按大背景块与板块标题位置推断"


def bounding_rect(rects: Iterable[Rect]) -> Rect:
    rects = list(rects)
    left = min(r.left for r in rects)
    top = min(r.top for r in rects)
    right = max(r.right for r in rects)
    bottom = max(r.bottom for r in rects)
    return Rect(left, top, right - left, bottom - top)


def pictures_in_region(slide, region: Rect):
    pics = []
    for pic in iter_picture_shapes(slide):
        pr = rect_of(pic)
        if pr.area and overlap_area(pr, region) / pr.area >= 0.45:
            pics.append(pic)
    return pics


def picture_in_top_right(panel: Rect, pic) -> tuple[bool, str]:
    """判断单张图片是否位于板块右上"角"。

    仅"中心在右半上半区"不足以判定为"右上角"（可能贴中线或偏中央）。
    这里同时校验：
    - 图片中心在板块右半区且上半区（原有近似判定）；
    - 图片右边界靠近板块右侧（距离 ≤ 板块宽的 20%）；
    - 图片上边界靠近板块顶部（距离 ≤ 板块高的 25%）；
    - 图片尺寸不至于占满整个板块（面积占板块 ≤ 55%），确保是"角落"而非"覆盖"。
    """
    pr = rect_of(pic)
    center_ok = pr.cx >= panel.left + panel.width * 0.58 and pr.cy <= panel.top + panel.height * 0.48

    right_gap = panel.right - pr.right
    top_gap = pr.top - panel.top
    right_close = 0 <= right_gap <= panel.width * 0.20
    top_close = 0 <= top_gap <= panel.height * 0.25

    area_ratio = pr.area / max(1, panel.area)
    not_covering = area_ratio <= 0.55

    ok = center_ok and right_close and top_close and not_covering
    detail = (
        f"中心右上={center_ok}，右缘距={right_gap}(≤{int(panel.width * 0.20)})={right_close}，"
        f"上缘距={top_gap}(≤{int(panel.height * 0.25)})={top_close}，"
        f"占板块面积={area_ratio:.2f}(≤0.55)={not_covering}"
    )
    return ok, detail


def _iter_potential_content_shapes(slide):
    """遍历幻灯片上"可能算作已有内容"的形状（文本框、自选图形、图表、表格等）。

    图片自身单独处理；纯装饰底板（尺寸接近整页/半页的大块填充色矩形）不视为内容。
    """
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        yield shape


def picture_top_right_is_blank(
    slide, panel: Rect, pic, other_pics_in_panel: list[object]
) -> tuple[bool, str]:
    """判断图片所在的板块右上角是否为"空白处"（不遮挡文字/其他内容）。

    "空白"按几何近似：图片矩形与板块内其它非图片形状（文本/自选图形/表格等）
    的重叠面积不超过图片面积的阈值。同时排除：
    - 图片自身；
    - 板块背景色块（面积 ≥ 板块面积 90% 的填充色形状）；
    - 完全包住图片的容器形状（如整块板块背景，图片本就该"放置"在其上）。
    """
    pr = rect_of(pic)
    slide_pic_iid = id(pic)
    overlaps: list[tuple[str, int]] = []
    total_overlap_ratio = 0.0
    for shape in _iter_potential_content_shapes(slide):
        sr = rect_of(shape)
        if sr.area == 0:
            continue
        # 只关心与图片有实际重叠、并且这个形状确实"落在"当前板块内的内容。
        if overlap_area(sr, panel) / max(1, sr.area) < 0.30:
            continue
        # 板块背景色块（尺寸接近整个板块的大填充块）不算内容遮挡。
        if sr.area >= panel.area * 0.90 and overlap_area(sr, panel) / max(1, sr.area) >= 0.85:
            continue
        ov = overlap_area(sr, pr)
        if ov <= 0:
            continue
        # 完全包住图片的容器（比如整块板块底板），不算遮挡。
        if ov >= pr.area * 0.98 and sr.area >= pr.area * 2:
            continue
        text = shape_text(shape).strip()
        label = text[:20] if text else f"<{shape.shape_type}>"
        overlaps.append((label, ov))
        total_overlap_ratio += ov / max(1, pr.area)

    # 也把板块内的"其它图片"视为遮挡（严格意义上此细则要求本板块仅有一张图片）。
    for other in other_pics_in_panel:
        if id(other) == slide_pic_iid:
            continue
        ov = overlap_area(rect_of(other), pr)
        if ov > 0:
            other_name = getattr(other, "name", "?")
            overlaps.append((f"<其他图片:{other_name}>", ov))
            total_overlap_ratio += ov / max(1, pr.area)

    is_blank = total_overlap_ratio <= 0.05  # 允许 ≤5% 的边缘擦碰
    if overlaps:
        summary = "、".join(f"{lbl}({ov})" for lbl, ov in overlaps[:3])
        detail = f"重叠内容={summary}，累计遮挡={total_overlap_ratio:.2%}"
    else:
        detail = "板块内右上区域无内容遮挡"
    return is_blank, detail


def evaluate_dimension1(path: Path) -> tuple[bool, str, Optional[Presentation]]:
    """维度1：交付文件为 .pptx 格式，且可正常打开。

    仅"能被 python-pptx 解析"不足以证明可正常打开，此处补充以下检查：
    - 扩展名与文件头一致（.pptx 必须是 ZIP，不能是加密 .pptx 的 OLE 复合文档头）；
    - ZIP 内部结构完整、必需部件齐全（否则 PowerPoint 打开会报错）；
    - 未设置打开密码（否则无法打开）。
    """
    # 1) 扩展名
    if path.suffix.lower() != ".pptx":
        return False, f"交付文件扩展名为 {path.suffix!r}，不是 .pptx", None
    if not path.exists():
        return False, f"文件不存在：{path}", None

    # 2) 文件头识别：加密 .pptx 实际为 OLE 复合文档，无法直接编辑/播放
    try:
        with open(path, "rb") as fh:
            header = fh.read(8)
    except OSError as exc:
        return False, f"文件读取失败，无法打开：{exc}", None
    if header.startswith(b"\xD0\xCF\x11\xE0\xA1\xB1\x1A\xE1"):
        return False, "文件为加密 OLE 复合文档（疑似设置了打开密码），无法正常编辑/播放", None
    if not header.startswith(b"PK"):
        return False, "文件不是有效的 .pptx（缺少 ZIP 文件头）", None

    # 3) ZIP 结构完整性 + 必需部件 + 保护属性检查
    pres_xml: bytes = b""
    try:
        with zipfile.ZipFile(str(path)) as zf:
            broken = zf.testzip()
            if broken is not None:
                return False, f"文件内部条目 {broken} 损坏，PowerPoint 打开会报错", None
            names = set(zf.namelist())
            required_parts = ["[Content_Types].xml", "ppt/presentation.xml"]
            missing = [p for p in required_parts if p not in names]
            if missing:
                return False, f"缺少必需部件 {missing}，无法正常打开/编辑", None
            pres_xml = zf.read("ppt/presentation.xml")
    except zipfile.BadZipFile:
        return False, "文件不是有效的 .pptx（ZIP 结构损坏）", None
    except Exception as exc:  # noqa: BLE001
        return False, f"文件结构检查失败：{exc}", None

    try:
        etree.fromstring(pres_xml)  # type: ignore[arg-type]
    except Exception as exc:  # noqa: BLE001
        return False, f"presentation.xml 解析失败：{exc}", None

    # 4) python-pptx 解析（可正常打开的基本前提：XML 语义合法）
    try:
        prs = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        return False, f"文件无法正常打开解析：{exc}", None

    slide_count = len(prs.slides)
    if slide_count == 0:
        return False, "文件可打开但没有幻灯片", None

    detail = (
        f"文件为 .pptx；ZIP 结构完整、无加密；共 {slide_count} 页；"
        f"判定为可正常打开"
    )
    return True, detail, prs


def check_title_wordart(prs: Presentation, path: Path) -> CheckResult:
    name = "首页标题“田间共生观察”改为白色外边蓝色填充的艺术字"
    if len(prs.slides) < 1:
        return CheckResult(5, name, False, "缺少首页")
    root = get_slide_xml_root(path, 1)
    theme_colors = _load_theme_colors(path)
    nodes = text_shape_xml_nodes_with_text(root, "田间共生观察")
    for sp in nodes:
        ok, why = has_wordart_white_line_blue_fill(sp, theme_colors)
        if ok:
            return CheckResult(5, name, True, why)
    return CheckResult(
        5, name, False,
        "未同时检测到蓝色文字填充、白色文字轮廓与艺术字特征（普通文本框着色不视为艺术字）",
    )


def check_slide2_panels_1_to_1(prs: Presentation) -> CheckResult:
    name = "第二页“田间防护：护鸟带与驱鸟器应用说明”下方护鸟带和驱鸟器板块分成左右两部分比例为1:1"
    if len(prs.slides) < 2:
        return CheckResult(5, name, False, "缺少第二页")
    slide = prs.slides[1]
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    panels = infer_content_panels(slide, slide_w, slide_h)
    if panels is None:
        return CheckResult(5, name, False, "未找到护鸟带/驱鸟器板块")
    bird_panel, repeller_panel, reason = panels

    # 细则点1：位于“田间防护：护鸟带与驱鸟器应用说明”标题下方。
    title_rect = find_slide2_main_title_rect(slide)
    if title_rect is None:
        below_title = False
        title_detail = "未找到指定标题"
    else:
        below_title = bird_panel.top >= title_rect.bottom and repeller_panel.top >= title_rect.bottom
        title_detail = f"标题底部={title_rect.bottom}，护鸟带顶={bird_panel.top}，驱鸟器顶={repeller_panel.top}"

    # 细则点2：护鸟带在左、驱鸟器在右（左右顺序固定，且中心 x 有明显分离）。
    order_ok = bird_panel.cx < repeller_panel.cx
    horizontal_sep = abs(bird_panel.cx - repeller_panel.cx) > slide_w * 0.20
    left_right = order_ok and horizontal_sep

    # 细则点3：两板块位于同一行（顶部/底部/中心 y 接近，且垂直上重叠明显）。
    max_h = max(bird_panel.height, repeller_panel.height, 1)
    same_row_center = abs(bird_panel.cy - repeller_panel.cy) <= max_h * 0.20
    v_overlap = max(
        0,
        min(bird_panel.bottom, repeller_panel.bottom) - max(bird_panel.top, repeller_panel.top),
    )
    same_row_overlap = v_overlap >= min(bird_panel.height, repeller_panel.height) * 0.60
    same_row = same_row_center and same_row_overlap

    # 细则点4：两板块"共同构成"标题下方的主要内容区。
    # 用两板块联合覆盖的宽度和高度分别占标题下方剩余画布的比例来衡量。
    if title_rect is not None:
        content_top = title_rect.bottom
    else:
        content_top = min(bird_panel.top, repeller_panel.top)
    content_h = max(1, slide_h - content_top)
    union_left = min(bird_panel.left, repeller_panel.left)
    union_right = max(bird_panel.right, repeller_panel.right)
    union_top = min(bird_panel.top, repeller_panel.top)
    union_bottom = max(bird_panel.bottom, repeller_panel.bottom)
    width_cover = (union_right - union_left) / max(1, slide_w)
    height_cover = (union_bottom - union_top) / content_h
    covers_content_area = width_cover >= 0.70 and height_cover >= 0.50

    # 细则点5：宽度与面积均满足 1:1（取更严格者，容差 ±8%）。
    balanced_width = two_sides_balanced(bird_panel, repeller_panel)
    if min(bird_panel.area, repeller_panel.area) > 0:
        area_ratio = bird_panel.area / repeller_panel.area
        balanced_area = abs(area_ratio - 1.0) <= 0.10
    else:
        area_ratio = 0.0
        balanced_area = False
    balanced = balanced_width and balanced_area

    ok = below_title and left_right and same_row and covers_content_area and balanced
    detail = (
        f"{reason}；{title_detail}；"
        f"护鸟带在左={order_ok}，水平分离={horizontal_sep}；"
        f"同一行(中心y≈)={same_row_center}，垂直重叠={same_row_overlap}；"
        f"覆盖标题下方内容区(宽{width_cover:.2f}/高{height_cover:.2f})={covers_content_area}；"
        f"护鸟带宽={bird_panel.width}、驱鸟器宽={repeller_panel.width}，宽度1:1={balanced_width}；"
        f"面积比={area_ratio:.3f}，面积1:1={balanced_area}"
    )
    return CheckResult(5, name, ok, detail)


def check_slide2_panel_images_top_right(prs: Presentation) -> CheckResult:
    name = "第二页护鸟带和驱鸟器板块分别有一张图片需放置在各自板块的右上角空白处"
    if len(prs.slides) < 2:
        return CheckResult(5, name, False, "缺少第二页")
    slide = prs.slides[1]
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    panels = infer_content_panels(slide, slide_w, slide_h)
    if panels is None:
        return CheckResult(5, name, False, "未找到护鸟带/驱鸟器板块")
    bird_panel, repeller_panel, _ = panels
    bird_pics = pictures_in_region(slide, bird_panel)
    repeller_pics = pictures_in_region(slide, repeller_panel)

    def _judge(panel: Rect, pics: list, label: str) -> tuple[bool, str]:  # type: ignore[type-arg]
        if len(pics) != 1:
            return False, f"{label}板块图片{len(pics)}张(需恰好1张)"
        pic = pics[0]
        corner_ok, corner_detail = picture_in_top_right(panel, pic)
        blank_ok, blank_detail = picture_top_right_is_blank(slide, panel, pic, pics)
        return corner_ok and blank_ok, f"{label}：{corner_detail}；{blank_detail}"

    ok_bird, detail_bird = _judge(bird_panel, bird_pics, "护鸟带")
    ok_repeller, detail_repeller = _judge(repeller_panel, repeller_pics, "驱鸟器")

    ok = ok_bird and ok_repeller
    detail = f"{detail_bird}；{detail_repeller}"
    return CheckResult(5, name, ok, detail)


def check_slide2_picture_not_square(prs: Presentation) -> CheckResult:
    name = "第二页护鸟带和驱鸟器图片横纵比不满足1:1"
    if len(prs.slides) < 2:
        return CheckResult(-5, name, False, "缺少第二页")
    slide = prs.slides[1]
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    panels = infer_content_panels(slide, slide_w, slide_h)
    if panels is None:
        return CheckResult(-5, name, False, "未找到护鸟带/驱鸟器板块")
    bird_panel, repeller_panel, _ = panels

    # 细则只针对护鸟带和驱鸟器两个板块内的图片。
    panel_pics = pictures_in_region(slide, bird_panel) + pictures_in_region(slide, repeller_panel)
    if not panel_pics:
        return CheckResult(-5, name, False, "护鸟带/驱鸟器板块内未检测到图片")

    # 细则点：图片横纵比不满足1:1（即存在任一图片不是1:1）时扣分。
    not_square = [p for p in panel_pics if not is_square_picture(p)]
    detail = "；".join(f"{p.name}: {picture_aspect_ratio(p):.3f}" for p in panel_pics)
    return CheckResult(-5, name, bool(not_square), detail)


def _slide_background_fill_hex(pptx_path: Path, slide_number: int) -> Optional[str]:
    """从 slideN.xml 读取幻灯片级背景填充色。

    仅识别 ``<p:cSld>/<p:bg>`` 下的 ``a:solidFill``（srgbClr / sysClr）；渐变/图片背景返回 None。
    幻灯片自身未定义背景时不回溯 layout/master —— 那些颜色对"左右半部分不同色"的判定并不适用。
    """
    try:
        root = get_slide_xml_root(pptx_path, slide_number)
    except Exception:
        return None
    for solid in root.xpath(".//p:cSld/p:bg//a:solidFill", namespaces=NS):
        for child in solid:
            tag = etree.QName(child.tag).localname
            if tag == "srgbClr":
                val = child.get("val")
                if val:
                    return val.upper()
            if tag == "sysClr":
                val = child.get("lastClr")
                if val:
                    return val.upper()
    return None


def check_slide2_background_bad(prs: Presentation, path: Path) -> CheckResult:
    name = "第二页PPT背景填充不满足：左半部分护鸟带背景填充色为深绿色，右半部分填充色为浅绿色"
    if len(prs.slides) < 2:
        return CheckResult(-5, name, False, "缺少第二页")
    slide = prs.slides[1]
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    slide_area = slide_w * slide_h
    left_half = Rect(0, 0, slide_w // 2, slide_h)
    right_half = Rect(slide_w // 2, 0, slide_w - slide_w // 2, slide_h)

    # 1) 识别"真正的背景层"：幻灯片自身背景 + 大面积填充色块（跨越整张画布或整个半区）。
    #    小于该阈值的填充块视为装饰/内容块，不参与背景色判定。
    slide_bg_hex = _slide_background_fill_hex(path, 2)

    background_shapes: list[tuple[Rect, str]] = []
    for shape in slide.shapes:
        color = shape_fill_hex(shape)
        if not color:
            continue
        sr = rect_of(shape)
        if sr.area == 0:
            continue
        # 背景形状必须"足够大"：占整页 ≥ 40%，或几乎覆盖某个半区（≥ 85%）。
        covers_slide = sr.area / max(1, slide_area) >= 0.40
        covers_left_half = overlap_area(sr, left_half) / max(1, left_half.area) >= 0.85
        covers_right_half = overlap_area(sr, right_half) / max(1, right_half.area) >= 0.85
        if covers_slide or covers_left_half or covers_right_half:
            background_shapes.append((sr, color))

    def dominant_bg_color(half: Rect) -> tuple[Optional[str], float, list[tuple[str, float]]]:
        """返回该半区覆盖度最高的背景色、覆盖比例、以及全部候选(用于诊断)。"""
        # 幻灯片背景视为一个"铺满整张画布"的候选，覆盖比例 = 100%。
        coverage: dict[str, int] = {}
        if slide_bg_hex:
            coverage[slide_bg_hex] = coverage.get(slide_bg_hex, 0) + half.area
        for rect, color in background_shapes:
            ov = overlap_area(rect, half)
            if ov <= 0:
                continue
            coverage[color] = coverage.get(color, 0) + ov
        if not coverage:
            return None, 0.0, []
        ranked = sorted(coverage.items(), key=lambda kv: kv[1], reverse=True)
        top_color, top_ov = ranked[0]
        top_ratio = top_ov / max(1, half.area)
        diagnostic = [(c, ov / max(1, half.area)) for c, ov in ranked[:3]]
        return top_color, top_ratio, diagnostic

    left_color, left_ratio, left_diag = dominant_bg_color(left_half)
    right_color, right_ratio, right_diag = dominant_bg_color(right_half)

    # 2) 要求主导背景色确实占半区的主要面积（≥ 60%），才判定其为"半区背景填充色"。
    left_dominant_enough = left_ratio >= 0.60
    right_dominant_enough = right_ratio >= 0.60

    left_is_deep = left_color is not None and classify_green(left_color) == "deep_green"
    right_is_light = right_color is not None and classify_green(right_color) == "light_green"

    ok_positive = left_dominant_enough and right_dominant_enough and left_is_deep and right_is_light

    def _fmt(diag: list[tuple[str, float]]) -> str:
        return "、".join(f"{c}({r:.0%})" for c, r in diag) if diag else "无"

    detail = (
        f"幻灯片背景={slide_bg_hex or '无'}；"
        f"左半区主导={left_color}({left_ratio:.0%})[{_fmt(left_diag)}]，深绿色={left_is_deep}(主导≥60%={left_dominant_enough})；"
        f"右半区主导={right_color}({right_ratio:.0%})[{_fmt(right_diag)}]，浅绿色={right_is_light}(主导≥60%={right_dominant_enough})"
    )
    # 扣分项：不满足即命中扣分。
    return CheckResult(-5, name, not ok_positive, detail)


def check_title_font_not_60(prs: Presentation, path: Path) -> CheckResult:
    name = "首页标题“田间共生观察”出现字号不是60磅"
    if len(prs.slides) < 1:
        return CheckResult(-1, name, False, "缺少首页")

    # 按 OOXML 继承链解析：run.rPr → 段落 pPr → 形状 lstStyle → slideLayout → slideMaster。
    entries = collect_title_effective_sizes(path, 1, "田间共生观察")
    if not entries:
        # 找不到含标题文本的 run。细则语义偏向"标题异常"，从严扣分。
        return CheckResult(-1, name, True, "未在首页 slide/layout/master 中定位到标题 run")

    resolved = [(s, note) for s, note in entries if s > 0]
    unresolved = [note for s, note in entries if s <= 0]

    bad_sizes = [s for s, _ in resolved if abs(s - 60) > 0.5]
    detail_parts = [note for _, note in entries]

    # 细则："出现字号不是 60 磅"即扣分。
    # - 存在非 60 磅的 run  →  命中扣分。
    # - 至少一个 run 未解析出字号  →  从严扣分（不再默认放过），并说明原因。
    if bad_sizes:
        hit = True
        reason = f"检测到非60磅字号={bad_sizes}"
    elif unresolved:
        hit = True
        reason = f"存在无法解析的字号，从严判定为异常：{unresolved}"
    else:
        hit = False
        reason = "所有 run 的最终字号均为 60 磅"

    return CheckResult(-1, name, hit, f"{reason}；明细：{'；'.join(detail_parts)}")


def check_slide2_title_missing(prs: Presentation) -> CheckResult:
    name = "第二页PPT顶部未出现“田间防护：护鸟带与驱鸟器应用说明”文字内容"
    expected = "田间防护：护鸟带与驱鸟器应用说明"
    if len(prs.slides) < 2:
        return CheckResult(-3, name, True, "缺少第二页")
    slide = prs.slides[1]
    slide_h = int(prs.slide_height)

    # 细则点：第二页“顶部”出现指定文字内容。找到含该文字的形状并判断其是否位于顶部。
    target = normalized_text(expected)
    title_shape = None
    for shape in slide.shapes:
        if target in normalized_text(shape_text(shape)):
            title_shape = shape
            break

    if title_shape is None:
        return CheckResult(-3, name, True, "未检测到指定文字内容")

    at_top = rect_of(title_shape).cy <= slide_h * 0.30
    detail = f"检测到文字内容，位于顶部={at_top}（中心y={int(rect_of(title_shape).cy)}）"
    # 未出现在顶部即命中扣分。
    return CheckResult(-3, name, not at_top, detail)


def check_bird_band_content_missing(prs: Presentation) -> CheckResult:
    name = "第二页护鸟带板块“材质”“原理”“设计”三条内容缺少任意一条"
    required = ["材质", "原理", "设计"]
    if len(prs.slides) < 2:
        return CheckResult(-3, name, True, "缺少第二页")
    slide = prs.slides[1]
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    # 细则点：范围限定在“护鸟带板块”内。先定位护鸟带板块区域，再取该区域内文字。
    panels = infer_content_panels(slide, slide_w, slide_h)
    if panels is None:
        return CheckResult(-3, name, True, "未找到护鸟带板块")
    bird_panel, _, _ = panels

    texts = []
    for shape in slide.shapes:
        t = shape_text(shape)
        if not t:
            continue
        sr = rect_of(shape)
        if sr.area and overlap_area(sr, bird_panel) / sr.area >= 0.5:
            texts.append(t)
    bird_text = "\n".join(texts)

    # 细则点：“材质”“原理”“设计”三条缺少任意一条即扣分。
    missing = [word for word in required if word not in bird_text]
    return CheckResult(-3, name, bool(missing), f"缺少：{missing}" if missing else "三条内容均检测到")


SCRIPT_ID = "071"
SUPPORTED_SUFFIXES = {".pptx"}


def _locate_target_file(dir_path: Path) -> Path | None:
    """在给定目录中寻找待评估的 .pptx 文件。

    - 忽略脚本自身、缓存目录以及以 ``~$`` 开头的 Office 临时文件；
    - 只识别 ``.pptx``；若存在多个候选，取修改时间最新的一个。
    """
    if not dir_path.is_dir():
        return None
    candidates: list[Path] = []
    for entry in dir_path.iterdir():
        if not entry.is_file():
            continue
        if entry.name.startswith("~$"):
            continue
        if entry.suffix.lower() not in SUPPORTED_SUFFIXES:
            continue
        candidates.append(entry)
    if not candidates:
        return None
    candidates.sort(key=lambda p: -p.stat().st_mtime)
    return candidates[0]


def evaluate(dir_path: str) -> dict[str, object]:
    """评估脚本所在目录里的 PPT 文件，返回结构化结果字典。

    参数 ``dir_path`` 为“脚本所在目录的路径”，脚本自行在该目录中查找并打开
    被评估的文档。返回结构遵循《脚本接口差异与统一建议.md》§2.2。
    """
    result: dict[str, object] = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": 0,
    }

    # 维度二各评分项的满分之和（正分累加，扣分项以其绝对值不计入满分基准，
    # 这里保留“正向满分”语义：10 项细则中 3 项为得分项、7 项为扣分项，
    # 满分为得分项之和 15 分）。
    max_positive_score = 5 + 5 + 5
    result["max_score"] = max_positive_score

    try:
        directory = Path(dir_path)
        target = _locate_target_file(directory)
        if target is None:
            result["status"] = "error"
            result["error"] = f"目录中未找到 .pptx 文件：{directory}"
            return result

        result["file_name"] = target.name

        dim1_ok, dim1_detail, prs = evaluate_dimension1(target)
        result["dim1_pass"] = dim1_ok
        if not dim1_ok or prs is None:
            result["dim1_reason"] = dim1_detail
            result["total_score"] = 0
            return result

        checks: list[CheckResult] = [
            check_title_wordart(prs, target),
            check_slide2_panels_1_to_1(prs),
            check_slide2_panel_images_top_right(prs),
            check_slide2_picture_not_square(prs),
            check_slide2_background_bad(prs, target),
            check_title_font_not_60(prs, target),
            check_slide2_title_missing(prs),
            check_bird_band_content_missing(prs),
        ]

        total = 0
        items: list[dict[str, object]] = []
        for item in checks:
            delta = item.score if item.passed else 0
            total += delta
            items.append(
                {
                    "rule": item.name,
                    "max_delta": item.score,
                    "delta": delta,
                    "hit": item.passed,
                    "detail": "",
                }
            )
        result["dim2_items"] = items
        result["total_score"] = total
        return result
    except Exception as exc:  # noqa: BLE001
        result["status"] = "error"
        result["error"] = f"{type(exc).__name__}: {exc}"
        return result


def main(argv: list[str]) -> int:
    """本地调试入口：把评估结果以 JSON 形式打印到标准输出。"""
    target_dir = argv[0] if argv else str(Path(__file__).resolve().parent)
    report = evaluate(target_dir)
    text = json.dumps(report, ensure_ascii=False, indent=2)
    try:
        print(text)
    except UnicodeEncodeError:
        # 兼容 Windows 默认 cp1252 控制台
        sys.stdout.buffer.write(text.encode("utf-8", errors="replace"))
        sys.stdout.buffer.write(b"\n")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
