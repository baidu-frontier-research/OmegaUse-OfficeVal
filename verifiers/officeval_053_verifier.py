#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""自动评估《电化学分析示意图》可编辑 PPT。

评估逻辑严格按题目要求分两层：
1. 维度1为硬门槛，任一项不满足则直接 0 分并跳过维度2。
2. 维度1通过后，维度2逐项检测；命中任一得分/扣分点即累计对应分值。

用法：
    python evaluate_ppt.py
    python evaluate_ppt.py "电化学分析示意图_可编辑PPT.pptx"
"""

from __future__ import annotations

import json
import math
import re
import sys
import xml.etree.ElementTree as ET
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable, Optional

from pptx import Presentation

EMU_PER_CM = 360_000
EMU_PER_PT = 12_700
DEFAULT_FILE = "电化学分析示意图_可编辑PPT.pptx"

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

SANS_FONTS = (
    "微软雅黑",
    "microsoft yahei",
    "黑体",
    "simhei",
    "heiti",
    "等线",
    "dengxian",
    "arial",
    "calibri",
    "noto sans",
    "sans",
)

# 直线/连接符类几何：line 是自由绘制的直线，straightConnector1 / bentConnectorN /
# curvedConnectorN 是从形状端点拉出的连接符。所有这些都归为"线形"对象，
# 只要满足箭头/线宽/颜色等要求都可作为箭头或连接线看待。
LINE_LIKE_GEOMS = {
    "line",
    "straightConnector1",
    "bentConnector2", "bentConnector3", "bentConnector4", "bentConnector5",
    "curvedConnector2", "curvedConnector3", "curvedConnector4", "curvedConnector5",
}

# 预设为箭头的形状（Office "块箭头" 分类）：本身即为一体箭头，
# 无需再依赖 headEnd/tailEnd 判定。
ARROW_PRESET_GEOMS = {
    "rightArrow", "leftArrow", "upArrow", "downArrow",
    "leftRightArrow", "upDownArrow", "quadArrow", "leftRightUpArrow",
    "bentArrow", "uturnArrow", "circularArrow",
    "curvedRightArrow", "curvedLeftArrow", "curvedUpArrow", "curvedDownArrow",
    "stripedRightArrow", "notchedRightArrow", "homePlate", "chevron",
    "bentUpArrow", "leftUpArrow",
}


@dataclass
class Box:
    left: float
    top: float
    width: float
    height: float

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height

    @property
    def x1(self) -> float:
        return min(self.left, self.right)

    @property
    def x2(self) -> float:
        return max(self.left, self.right)

    @property
    def y1(self) -> float:
        return min(self.top, self.bottom)

    @property
    def y2(self) -> float:
        return max(self.top, self.bottom)

    @property
    def cx(self) -> float:
        return (self.x1 + self.x2) / 2

    @property
    def cy(self) -> float:
        return (self.y1 + self.y2) / 2

    @property
    def area(self) -> float:
        return max(0.0, self.x2 - self.x1) * max(0.0, self.y2 - self.y1)


@dataclass
class ShapeInfo:
    index: int
    name: str
    geom: str
    box: Box
    text: str
    fill_colors: list[tuple[int, int, int]]
    line_color: Optional[tuple[int, int, int]]
    line_width: Optional[float]
    dashed: bool
    arrow: bool
    font_name: Optional[str]
    font_size: Optional[float]
    text_color: Optional[tuple[int, int, int]]
    text_align: Optional[str]
    has_blip: bool
    has_pattern_fill: bool
    has_transparency: bool
    # 形状本身的 a:xfrm 旋转(度)。PowerPoint 约定: 正值=顺时针(屏幕上 Y 向下),
    # 未设置时为 0.0; 用于判定"向左倾斜"等方向性 rubric。
    rotation: float = 0.0
    flip_h: bool = False
    flip_v: bool = False

    @property
    def x1(self) -> float:
        return self.box.x1

    @property
    def x2(self) -> float:
        return self.box.x2

    @property
    def y1(self) -> float:
        return self.box.y1

    @property
    def y2(self) -> float:
        return self.box.y2

    @property
    def cx(self) -> float:
        return self.box.cx

    @property
    def cy(self) -> float:
        return self.box.cy

    @property
    def w_abs(self) -> float:
        return self.box.x2 - self.box.x1

    @property
    def h_abs(self) -> float:
        return self.box.y2 - self.box.y1

    @property
    def length(self) -> float:
        return math.hypot(self.box.width, self.box.height)

    @property
    def angle(self) -> float:
        # PPT 坐标向下为正，因此这里使用原始 dy 判断走向即可。
        return math.degrees(math.atan2(self.box.height, self.box.width))

    @property
    def is_text(self) -> bool:
        return bool(norm_text(self.text))

    @property
    def is_line(self) -> bool:
        # line、straightConnector 与其他 connector 都视为线；同时兼容手绘的极细矩形。
        return (
            self.geom in LINE_LIKE_GEOMS
            or (self.length > 0.2 and (self.w_abs < 0.08 or self.h_abs < 0.08))
        )

    @property
    def is_arrow_shape(self) -> bool:
        """一体箭头：预设箭头几何 或 line/连接符设置了 head/tail 箭头。"""
        if self.geom in ARROW_PRESET_GEOMS:
            return True
        return self.arrow and self.geom in LINE_LIKE_GEOMS

    @property
    def is_rect_like(self) -> bool:
        return self.geom in {"rect", "roundRect", "parallelogram", "trapezoid"}

    @property
    def is_ellipse(self) -> bool:
        return self.geom == "ellipse"


@dataclass
class EvaluationContext:
    path: Path
    prs: object
    slide_width: float
    slide_height: float
    slide_count: int
    shapes: list[ShapeInfo]
    open_error: Optional[str] = None

    @property
    def text_shapes(self) -> list[ShapeInfo]:
        return [s for s in self.shapes if s.is_text]

    @property
    def line_shapes(self) -> list[ShapeInfo]:
        return [s for s in self.shapes if s.is_line]

    @property
    def rect_shapes(self) -> list[ShapeInfo]:
        return [s for s in self.shapes if s.is_rect_like]

    @property
    def ellipse_shapes(self) -> list[ShapeInfo]:
        return [s for s in self.shapes if s.is_ellipse]


@dataclass
class RuleResult:
    label: str
    hit: bool
    detail: str
    points: int = 0
    # 仅当该规则在执行时抛出异常才为 True；用于在输出中区分
    # "正常命中/未命中" 与 "异常" 两种状态。
    errored: bool = False


@dataclass
class ScoreItem:
    points: int
    label: str
    checker: Callable[[EvaluationContext], RuleResult]


def cm(value) -> float:
    return float(value or 0) / EMU_PER_CM


def pt(value) -> Optional[float]:
    if value is None:
        return None
    return float(value) / EMU_PER_PT


def norm_text(text: str) -> str:
    text = (text or "").replace(" ", " ")
    text = text.replace("（", "(").replace("）", ")")
    text = text.replace("／", "/").replace("—", "-").replace("–", "-")
    return re.sub(r"\s+", "", text).strip().lower()


def display_text(text: str) -> str:
    return re.sub(r"\s+", " ", (text or "").strip())


def rgb_from_hex(value: str | None) -> Optional[tuple[int, int, int]]:
    if not value:
        return None
    value = value.strip().lstrip("#")[-6:]
    if len(value) != 6:
        return None
    try:
        return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return None


def color_distance(a: Optional[tuple[int, int, int]], b: tuple[int, int, int]) -> float:
    if a is None:
        return 999.0
    return math.sqrt(sum((x - y) ** 2 for x, y in zip(a, b)))


def near_color(color: Optional[tuple[int, int, int]], target: tuple[int, int, int], tol: float = 80) -> bool:
    return color_distance(color, target) <= tol


def any_near(colors: Iterable[Optional[tuple[int, int, int]]], target: tuple[int, int, int], tol: float = 80) -> bool:
    return any(near_color(c, target, tol) for c in colors if c is not None)


def is_whiteish(color: Optional[tuple[int, int, int]]) -> bool:
    return color is not None and min(color) >= 230


def is_blackish(color: Optional[tuple[int, int, int]]) -> bool:
    return color is not None and max(color) <= 100


def is_dark(color: Optional[tuple[int, int, int]]) -> bool:
    return color is not None and sum(color) / 3 <= 110


def is_grayish(color: Optional[tuple[int, int, int]]) -> bool:
    # 判定为"灰色"只要 R/G/B 三通道差异较小即可，不再限定为主题灰的窄区间。
    return color is not None and max(color) - min(color) <= 60


def is_light_grayish(color: Optional[tuple[int, int, int]]) -> bool:
    # 只要属于灰色家族即视为合规，不再强制"浅灰"这一具体色调。
    return is_grayish(color)


def is_blue(color: Optional[tuple[int, int, int]]) -> bool:
    if color is None:
        return False
    r, g, b = color
    return b >= 120 and b >= r + 30 and b >= g - 20


def is_medium_blue(color: Optional[tuple[int, int, int]]) -> bool:
    # 只要属于蓝色家族即视为合规，不再要求"中蓝"这一具体色调。
    return is_blue(color)


def is_dark_blue(color: Optional[tuple[int, int, int]]) -> bool:
    # 只要属于蓝色家族即视为合规，不再要求"深蓝"这一具体色调。
    return is_blue(color)


def is_light_blue(color: Optional[tuple[int, int, int]]) -> bool:
    # 只要属于蓝色家族即视为合规，不再要求"浅蓝"这一具体色调。
    return is_blue(color)


def is_green(color: Optional[tuple[int, int, int]]) -> bool:
    if color is None:
        return False
    r, g, b = color
    return g >= 100 and g >= r + 20 and g >= b + 5


def is_red(color: Optional[tuple[int, int, int]]) -> bool:
    if color is None:
        return False
    r, g, b = color
    return r >= 130 and r >= g + 30 and r >= b + 30


def is_cyan_or_whiteblue(color: Optional[tuple[int, int, int]]) -> bool:
    # 青色/白蓝色只要属于蓝色家族或浅色/灰色即视为合规。
    if color is None:
        return False
    return is_blue(color) or is_whiteish(color) or is_grayish(color)


def in_range(value: float, low: float, high: float, tol: float = 0.0) -> bool:
    return low - tol <= value <= high + tol


# PPT 中若未显式设置线宽，默认按 1.0 磅处理（Office/PowerPoint 默认值）。
DEFAULT_LINE_WIDTH_PT = 1.0


def effective_line_width(width: Optional[float]) -> float:
    """线宽取值：未显式设置时按默认 1.0 磅处理。"""
    return DEFAULT_LINE_WIDTH_PT if width is None else width


def line_width_in_range(width: Optional[float], low: float, high: float, tol: float = 0.0) -> bool:
    """线宽范围判定：未显式设置线宽按默认 1.0 磅处理，只要 1.0 在范围内即视为通过。"""
    return in_range(effective_line_width(width), low, high, tol)


def center_in(s: ShapeInfo, left: float, top: float, right: float, bottom: float, tol: float = 0.25) -> bool:
    return in_range(s.cx, left, right, tol) and in_range(s.cy, top, bottom, tol)


def within_region(s: ShapeInfo, left: float, top: float, right: float, bottom: float, tol: float = 0.35) -> bool:
    return (
        s.x1 >= left - tol
        and s.y1 >= top - tol
        and s.x2 <= right + tol
        and s.y2 <= bottom + tol
    )


def intersects_region(s: ShapeInfo, left: float, top: float, right: float, bottom: float, tol: float = 0.0) -> bool:
    return not (s.x2 < left - tol or s.x1 > right + tol or s.y2 < top - tol or s.y1 > bottom + tol)


def overlap_area(a: Box, b: Box) -> float:
    w = max(0.0, min(a.x2, b.x2) - max(a.x1, b.x1))
    h = max(0.0, min(a.y2, b.y2) - max(a.y1, b.y1))
    return w * h


def union_bbox(shapes: Iterable[ShapeInfo]) -> Optional[Box]:
    vals = list(shapes)
    if not vals:
        return None
    x1 = min(s.x1 for s in vals)
    y1 = min(s.y1 for s in vals)
    x2 = max(s.x2 for s in vals)
    y2 = max(s.y2 for s in vals)
    return Box(x1, y1, x2 - x1, y2 - y1)


def first_xml_rgb(elem, theme: dict[str, tuple[int, int, int]] | None = None) -> tuple[int, int, int] | None:
    """返回 elem 内首个可解析的颜色 RGB。
    优先级：a:srgbClr → a:schemeClr(经 theme 反查) → None。
    theme: {schemeName: (r,g,b)} 由 _load_theme_colors 生成; 传 None 时仅解析 srgbClr。
    """
    for clr in elem.iter():
        tag = _local_tag(clr.tag)
        if tag == "srgbClr":
            rgb = rgb_from_hex(clr.get("val"))
            if rgb is not None:
                return _apply_shade_tint(rgb, clr)
        elif tag == "schemeClr" and theme is not None:
            name = clr.get("val")
            base = theme.get(name)
            if base is not None:
                return _apply_shade_tint(base, clr)
    return None


def all_xml_rgbs(elem, theme: dict[str, tuple[int, int, int]] | None = None) -> list[tuple[int, int, int]]:
    """返回 elem 内全部可解析颜色 RGB(srgbClr + schemeClr, 顺序保留)。"""
    colors: list[tuple[int, int, int]] = []
    for clr in elem.iter():
        tag = _local_tag(clr.tag)
        if tag == "srgbClr":
            rgb = rgb_from_hex(clr.get("val"))
            if rgb is not None:
                colors.append(_apply_shade_tint(rgb, clr))
        elif tag == "schemeClr" and theme is not None:
            name = clr.get("val")
            base = theme.get(name)
            if base is not None:
                colors.append(_apply_shade_tint(base, clr))
    return colors


def _local_tag(tag: str) -> str:
    """去掉命名空间前缀, 例如 '{http://.../main}srgbClr' → 'srgbClr'。"""
    if "}" in tag:
        return tag.split("}", 1)[1]
    return tag


def _apply_shade_tint(rgb: tuple[int, int, int], clr_elem) -> tuple[int, int, int]:
    """对 a:*Clr 子级的 a:lumMod / a:lumOff / a:shade / a:tint 做常见亮度调整。
    数值均以千分数表示; 未识别的调整忽略。目的是让主题色反查后得到更接近实际渲染的 RGB。
    """
    r, g, b = rgb
    for child in clr_elem:
        tag = _local_tag(child.tag)
        try:
            val = int(child.get("val", "0"))
        except (TypeError, ValueError):
            continue
        v = val / 100000.0
        if tag == "lumMod":
            r = int(r * v); g = int(g * v); b = int(b * v)
        elif tag == "lumOff":
            add = int(255 * v)
            r = min(255, r + add); g = min(255, g + add); b = min(255, b + add)
        elif tag == "shade":
            r = int(r * v); g = int(g * v); b = int(b * v)
        elif tag == "tint":
            # PowerPoint tint: mix toward white
            r = int(r + (255 - r) * (1 - v))
            g = int(g + (255 - g) * (1 - v))
            b = int(b + (255 - b) * (1 - v))
    return (max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b)))


# 主题 clrScheme 元素名 → dataclass 字段名(供 fallback / 别名映射用)
_SCHEME_ALIASES = {
    "bg1": "lt1",   # background 1 通常等价于 lt1
    "tx1": "dk1",   # text 1 通常等价于 dk1
    "bg2": "lt2",
    "tx2": "dk2",
}


def _load_theme_colors(pptx_path: Path) -> dict[str, tuple[int, int, int]]:
    """从 pptx 内的 ppt/theme/theme1.xml 里读 clrScheme, 返回 {schemeName: (r,g,b)}。
    包含 dk1/lt1/dk2/lt2/accent1-6/hlink/folHlink 以及 bg1/tx1/bg2/tx2 别名映射。
    读不到 / 解析失败时返回空 dict。
    """
    scheme: dict[str, tuple[int, int, int]] = {}
    try:
        with zipfile.ZipFile(str(pptx_path)) as zf:
            theme_name = None
            for n in zf.namelist():
                if n.startswith("ppt/theme/") and n.endswith(".xml"):
                    theme_name = n
                    break
            if theme_name is None:
                return scheme
            with zf.open(theme_name) as fh:
                data = fh.read()
        root = ET.fromstring(data)
        # 找 a:clrScheme
        for elem in root.iter():
            if _local_tag(elem.tag) == "clrScheme":
                for child in elem:
                    name = _local_tag(child.tag)  # dk1 / lt1 / accent1 / hlink / folHlink
                    # 每个 name 下应有一个 srgbClr 或 sysClr
                    for c in child.iter():
                        ctag = _local_tag(c.tag)
                        if ctag == "srgbClr":
                            rgb = rgb_from_hex(c.get("val"))
                            if rgb is not None:
                                scheme[name] = rgb
                                break
                        elif ctag == "sysClr":
                            # sysClr 通常带 lastClr 属性作为回退 RGB
                            last = c.get("lastClr")
                            rgb = rgb_from_hex(last) if last else None
                            if rgb is None:
                                # 常见 sysClr val: windowText → 黑, window → 白
                                sys_val = c.get("val", "").lower()
                                if sys_val in ("windowtext",):
                                    rgb = (0, 0, 0)
                                elif sys_val in ("window",):
                                    rgb = (255, 255, 255)
                            if rgb is not None:
                                scheme[name] = rgb
                                break
                break  # 只取第一个 clrScheme
    except (zipfile.BadZipFile, ET.ParseError, KeyError, OSError):
        return scheme
    # 别名兜底: bg1→lt1 等
    for alias, target in _SCHEME_ALIASES.items():
        if alias not in scheme and target in scheme:
            scheme[alias] = scheme[target]
        if target not in scheme and alias in scheme:
            scheme[target] = scheme[alias]
    return scheme


def get_prst_geom(elem) -> str:
    geom = elem.find(".//a:prstGeom", namespaces=NS)
    return geom.get("prst", "") if geom is not None else ""


def get_name(elem) -> str:
    c = elem.find(".//p:cNvPr", namespaces=NS)
    return c.get("name", "") if c is not None else ""


def get_text(elem) -> str:
    tx_body = elem.find(".//p:txBody", namespaces=NS)
    if tx_body is None:
        return ""
    paragraphs = []
    for p in tx_body.findall(".//a:p", namespaces=NS):
        chars = [t.text or "" for t in p.findall(".//a:t", namespaces=NS)]
        para = "".join(chars).strip()
        if para:
            paragraphs.append(para)
    return "\n".join(paragraphs)


def get_line_color(elem, theme: dict[str, tuple[int, int, int]] | None = None) -> tuple[int, int, int] | None:
    ln = elem.find(".//a:ln", namespaces=NS)
    if ln is None:
        return None
    return first_xml_rgb(ln, theme)


def get_line_width(elem) -> Optional[float]:
    ln = elem.find(".//a:ln", namespaces=NS)
    if ln is None:
        return None
    w = ln.get("w")
    if w is None:
        return None
    return pt(int(w))


def get_fill_colors(elem, theme: dict[str, tuple[int, int, int]] | None = None) -> list[tuple[int, int, int]]:
    sp_pr = elem.find(".//p:spPr", namespaces=NS)
    if sp_pr is None:
        return []
    fills = []
    for tag in ("solidFill", "gradFill"):
        for fill in sp_pr.findall(f".//a:{tag}", namespaces=NS):
            fills.extend(all_xml_rgbs(fill, theme))
    # Do not let line/text colors dominate fill; if no explicit fill, return [].
    return fills


def has_dash(elem) -> bool:
    ln = elem.find(".//a:ln", namespaces=NS)
    if ln is None:
        return False
    prst = ln.find("a:prstDash", namespaces=NS)
    if prst is not None and prst.get("val", "solid") not in {"solid", ""}:
        return True
    return ln.find("a:custDash", namespaces=NS) is not None


def has_arrow(elem) -> bool:
    for tag in ("headEnd", "tailEnd"):
        end = elem.find(f".//a:{tag}", namespaces=NS)
        if end is not None and end.get("type", "none") not in {"none", ""}:
            return True
    return False


def has_blip(elem) -> bool:
    return elem.find(".//a:blip", namespaces=NS) is not None


def has_pattern_fill(elem) -> bool:
    return elem.find(".//a:pattFill", namespaces=NS) is not None


def has_transparency(elem) -> bool:
    return elem.find(".//a:alpha", namespaces=NS) is not None


def get_xfrm_rotation(elem) -> tuple[float, bool, bool]:
    """返回形状 (rotation_deg, flipH, flipV)。
    从 p:spPr/a:xfrm 读 rot(EMU 单位 60000/度), flipH/flipV。未设置返回 (0.0, False, False)。
    """
    xfrm = elem.find(".//p:spPr/a:xfrm", namespaces=NS)
    if xfrm is None:
        # 部分连接符/组合的 xfrm 直接在 spPr 之外; 兜底再找一次
        xfrm = elem.find(".//a:xfrm", namespaces=NS)
    if xfrm is None:
        return 0.0, False, False
    rot_attr = xfrm.get("rot")
    try:
        rot = (int(rot_attr) / 60000.0) if rot_attr else 0.0
    except ValueError:
        rot = 0.0
    # 归一化到 (-180, 180]
    rot = ((rot + 180.0) % 360.0) - 180.0
    flip_h = xfrm.get("flipH") in ("1", "true")
    flip_v = xfrm.get("flipV") in ("1", "true")
    return rot, flip_h, flip_v


def slide_background_colors(slide, theme: dict[str, tuple[int, int, int]] | None = None) -> list[tuple[int, int, int]]:
    bg = slide._element.find(".//p:bg", namespaces=NS)
    if bg is None:
        return []
    return all_xml_rgbs(bg, theme)


def get_font_info(elem, theme: dict[str, tuple[int, int, int]] | None = None) -> tuple[Optional[str], Optional[float], Optional[tuple[int, int, int]]]:
    rpr = elem.find(".//a:rPr", namespaces=NS)
    if rpr is None:
        return None, None, None
    size = None
    if rpr.get("sz"):
        try:
            size = int(rpr.get("sz")) / 100.0
        except ValueError:
            size = None
    latin = rpr.find("a:latin", namespaces=NS)
    ea = rpr.find("a:ea", namespaces=NS)
    name = None
    if ea is not None and ea.get("typeface"):
        name = ea.get("typeface")
    elif latin is not None and latin.get("typeface"):
        name = latin.get("typeface")
    return name, size, first_xml_rgb(rpr, theme)


def get_text_align(elem) -> Optional[str]:
    """Return paragraph alignment: 'ctr', 'l', 'r', 'just', or None."""
    tx_body = elem.find(".//p:txBody", namespaces=NS)
    if tx_body is None:
        return None
    # Check body-level default (bodyPr / lstStyle) — paragraph overrides take precedence
    for p in tx_body.findall(".//a:p", namespaces=NS):
        ppr = p.find("a:pPr", namespaces=NS)
        if ppr is not None and ppr.get("algn"):
            return ppr.get("algn")
    return None


def build_context(path: Path) -> EvaluationContext:
    try:
        prs = Presentation(str(path))
        slide_count = len(prs.slides)
        slide_width = cm(prs.slide_width)
        slide_height = cm(prs.slide_height)
        # 一次性加载 theme1.xml 里的 clrScheme, 供后续所有颜色解析反查 schemeClr
        theme = _load_theme_colors(path)
        shapes: list[ShapeInfo] = []
        if slide_count:
            slide = prs.slides[0]
            for idx, raw in enumerate(slide.shapes, start=1):
                elem = raw._element
                b = Box(cm(raw.left), cm(raw.top), cm(raw.width), cm(raw.height))
                font_name, font_size, text_color = get_font_info(elem, theme)
                rot, flip_h, flip_v = get_xfrm_rotation(elem)
                shapes.append(
                    ShapeInfo(
                        index=idx,
                        name=get_name(elem),
                        geom=get_prst_geom(elem),
                        box=b,
                        text=get_text(elem),
                        fill_colors=get_fill_colors(elem, theme),
                        line_color=get_line_color(elem, theme),
                        line_width=get_line_width(elem),
                        dashed=has_dash(elem),
                        arrow=has_arrow(elem),
                        font_name=font_name,
                        font_size=font_size,
                        text_color=text_color,
                        text_align=get_text_align(elem),
                        has_blip=has_blip(elem),
                        has_pattern_fill=has_pattern_fill(elem),
                        has_transparency=has_transparency(elem),
                        rotation=rot,
                        flip_h=flip_h,
                        flip_v=flip_v,
                    )
                )
        return EvaluationContext(path, prs, slide_width, slide_height, slide_count, shapes)
    except Exception as exc:
        return EvaluationContext(path, None, 0, 0, 0, [], f"{type(exc).__name__}: {exc}")


def text_matches(s: ShapeInfo, expected: str | list[str] | tuple[str, ...]) -> bool:
    vals = expected if isinstance(expected, (list, tuple)) else [expected]
    n = norm_text(s.text)
    return any(norm_text(v) in n for v in vals)


def find_text(ctx: EvaluationContext, expected: str | list[str] | tuple[str, ...]) -> list[ShapeInfo]:
    return [s for s in ctx.text_shapes if text_matches(s, expected)]


def find_text_in(ctx: EvaluationContext, expected, left: float, top: float, right: float, bottom: float, tol: float = 0.45) -> list[ShapeInfo]:
    return [s for s in find_text(ctx, expected) if center_in(s, left, top, right, bottom, tol)]


def has_text_style(s: ShapeInfo, size_min: float, size_max: float, color_pred: Callable[[Optional[tuple[int, int, int]]], bool] | None = None, size_tol: float = 1.2) -> bool:
    size_ok = s.font_size is not None and in_range(s.font_size, size_min, size_max, size_tol)
    font_ok = True
    if s.font_name:
        font_ok = any(k in s.font_name.lower() for k in SANS_FONTS) or any(k in s.font_name for k in ("微软雅黑", "黑体", "等线"))
    color_ok = True if color_pred is None else (s.text_color is not None and color_pred(s.text_color))
    return size_ok and font_ok and color_ok


def count_shapes(region: tuple[float, float, float, float], shapes: Iterable[ShapeInfo], pred: Callable[[ShapeInfo], bool] = lambda s: True) -> int:
    l, t, r, b = region
    return sum(1 for s in shapes if pred(s) and intersects_region(s, l, t, r, b, 0.15))


def line_in_region(ctx: EvaluationContext, left: float, top: float, right: float, bottom: float, color_pred=None, width_range: Optional[tuple[float, float]] = None, dashed=None, arrow=None, horizontal=None, vertical=None, min_len: float = 0.0, max_len: Optional[float] = None, width_tol: float = 0.45, region_tol: float = 0.25, direction: Optional[str] = None) -> list[ShapeInfo]:
    vals = []
    # 当筛选"一体箭头"时，把预设箭头形状（rightArrow/downArrow 等 box 状的箭头）
    # 也纳入候选池；否则仅遍历 line/连接符 等线形对象。
    if arrow is True:
        candidates = [s for s in ctx.shapes if s.is_line or s.is_arrow_shape]
    else:
        candidates = ctx.line_shapes
    for s in candidates:
        if not intersects_region(s, left, top, right, bottom, region_tol):
            continue
        if color_pred and not color_pred(s.line_color):
            continue
        if width_range and not line_width_in_range(s.line_width, width_range[0], width_range[1], width_tol):
            continue
        if dashed is not None and s.dashed != dashed:
            continue
        if arrow is True:
            # 一体箭头：line/连接符带 head/tail，或者预设几何本身即为箭头形状
            if not s.is_arrow_shape:
                continue
        elif arrow is False:
            if s.is_arrow_shape:
                continue
        if horizontal is True and not (s.h_abs <= max(0.12, s.w_abs * 0.25)):
            continue
        if vertical is True and not (s.w_abs <= max(0.12, s.h_abs * 0.25)):
            continue
        if direction == "right" and s.box.width <= 0:
            continue
        if direction == "left" and s.box.width >= 0:
            continue
        if direction == "down" and s.box.height <= 0:
            continue
        if direction == "up" and s.box.height >= 0:
            continue
        if s.length < min_len:
            continue
        if max_len is not None and s.length > max_len:
            continue
        vals.append(s)
    return vals


def rect_in_region(ctx: EvaluationContext, left: float, top: float, right: float, bottom: float, *, dashed=None, color_pred=None, fill_pred=None, min_w=0.0, min_h=0.0) -> list[ShapeInfo]:
    vals = []
    for s in ctx.rect_shapes:
        if not center_in(s, left, top, right, bottom, 0.45):
            continue
        if dashed is not None and s.dashed != dashed:
            continue
        if color_pred and not color_pred(s.line_color):
            continue
        if fill_pred and not any(fill_pred(c) for c in s.fill_colors):
            continue
        if s.w_abs < min_w or s.h_abs < min_h:
            continue
        vals.append(s)
    return vals


def ellipses_in_region(ctx: EvaluationContext, left: float, top: float, right: float, bottom: float, *, fill_pred=None, line_pred=None, min_w=0.0, max_h=99.0) -> list[ShapeInfo]:
    vals = []
    for s in ctx.ellipse_shapes:
        if not center_in(s, left, top, right, bottom, 0.45):
            continue
        if fill_pred and not any(fill_pred(c) for c in s.fill_colors):
            continue
        if line_pred and not line_pred(s.line_color):
            continue
        if s.w_abs < min_w or s.h_abs > max_h:
            continue
        vals.append(s)
    return vals


def segmented_blue_border(ctx: EvaluationContext, left: float, top: float, right: float, bottom: float, min_segments: int = 16) -> tuple[bool, int]:
    """Detect editable dashed boxes drawn as many short blue line segments."""
    segments = line_in_region(ctx, left, top, right, bottom, color_pred=is_blue, width_range=(0.7, 1.8), min_len=0.08)
    short_segments = [s for s in segments if s.length <= 0.55]
    if not short_segments:
        return False, 0
    horizontals = [s for s in short_segments if s.h_abs <= max(0.08, s.w_abs * 0.35)]
    verticals = [s for s in short_segments if s.w_abs <= max(0.08, s.h_abs * 0.35)]
    x_span = max(s.x2 for s in short_segments) - min(s.x1 for s in short_segments)
    y_span = max(s.y2 for s in short_segments) - min(s.y1 for s in short_segments)
    expected_w = right - left
    expected_h = bottom - top
    spans_region = x_span >= expected_w * 0.90 and y_span >= expected_h * 0.80
    has_box_sides = len(horizontals) >= max(4, min_segments // 3) and len(verticals) >= 2
    hit = len(short_segments) >= min_segments and spans_region and has_box_sides
    return hit, len(short_segments)


def result(label: str, hit: bool, detail: str, points: int = 0) -> RuleResult:
    return RuleResult(label=label, hit=hit, detail=detail, points=points)


# ───────────────────────── 维度1 ─────────────────────────

def dimension1_checks(ctx: EvaluationContext) -> list[RuleResult]:
    checks: list[RuleResult] = []

    ext_ok = ctx.path.suffix.lower() == ".pptx"
    checks.append(result("交付文件为 .pptx 格式", ext_ok, f"扩展名：{ctx.path.suffix or '无'}"))

    open_ok = ctx.open_error is None and ctx.prs is not None
    checks.append(result("文件可正常打开", open_ok, "已成功读取 PPT" if open_ok else f"无法打开：{ctx.open_error}"))
    if not open_ok:
        return checks

    one_slide = ctx.slide_count == 1
    checks.append(result("交付PPT只包含1页幻灯片", one_slide, f"检测到 {ctx.slide_count} 页"))

    # 按用户指示已删除以下 rubric 项对应的检查(仅保留扩展名/文件可打开/幻灯片数量)：
    #   - 第1页所有可见元素均为可编辑对象(editable_count/text_count/picture_area 判定 → 已删除)
    #   - 第1页页面比例为横向16:9 / 无整体拉伸压缩旋转裁切(ratio + bbox 越界判定 → 已删除)
    #   - 无 >50% 空白 / 文字大面积重叠 / 对象越界(content_ratio + text_overlaps → 已删除)
    #   - "关键文本均为可编辑文本对象"锚点(14 组核心标签存在性 → 已删除)

    return checks


# ───────────────────────── 维度2规则 ─────────────────────────

def check_background(ctx: EvaluationContext) -> RuleResult:
    page_area = ctx.slide_width * ctx.slide_height
    full_page_shapes = [
        s for s in ctx.shapes
        if s.x1 <= 0.25 and s.y1 <= 0.25
        and s.x2 >= ctx.slide_width - 0.25 and s.y2 >= ctx.slide_height - 0.25
    ]
    full_page_rects = [s for s in full_page_shapes if s.is_rect_like]
    slide_bg_colors = slide_background_colors(ctx.prs.slides[0]) if ctx.prs is not None and ctx.slide_count else []

    full_page_picture_backgrounds = [s for s in full_page_shapes if s.has_blip]
    non_white_full_page_fills = [
        s for s in full_page_rects
        if s.fill_colors and not any(is_whiteish(c) for c in s.fill_colors)
    ]
    large_pattern_or_watermark_shapes = [
        s for s in ctx.shapes
        if s.box.area >= page_area * 0.55 and (s.has_pattern_fill or s.has_transparency)
    ]
    has_white_page_background = (
        not non_white_full_page_fills
        and (
            not slide_bg_colors
            or any(is_whiteish(c) for c in slide_bg_colors)
            or any(any(is_whiteish(c) for c in s.fill_colors) for s in full_page_rects)
        )
    )
    hit = has_white_page_background and not full_page_picture_backgrounds and not large_pattern_or_watermark_shapes
    return result(
        "第1页白色页面背景",
        hit,
        f"白/极浅灰白背景={'是' if has_white_page_background else '否'}，整页图片背景 {len(full_page_picture_backgrounds)} 个，大面积底纹/水印 {len(large_pattern_or_watermark_shapes)} 个",
    )


def check_outer_border(ctx: EvaluationContext) -> RuleResult:
    # rubric 修改后: 位于页面四周边缘, 蓝色单实线矩形, 线宽 0.5–3磅.
    # 相较原细则删除了"距页面四边 0–0.2cm 范围内"的硬约束,
    # 但仍需体现"位于页面四周边缘"—— 采用较宽松的边距容差 (≤1.0cm)
    # 以区分位于页面中部的普通蓝色矩形.
    EDGE_TOL = 1.0
    candidates = []
    for s in ctx.rect_shapes:
        if not s.is_rect_like:
            continue
        left_margin = s.x1
        top_margin = s.y1
        right_margin = ctx.slide_width - s.x2
        bottom_margin = ctx.slide_height - s.y2
        margins_ok = all(in_range(v, 0.0, EDGE_TOL, 0.0) for v in (left_margin, right_margin, top_margin, bottom_margin))
        line_ok = (
            not s.dashed
            and is_blue(s.line_color)
            and line_width_in_range(s.line_width, 0.5, 3.0, 0.0)
        )
        if margins_ok and line_ok:
            candidates.append(s)
    return result(
        "第1页外侧蓝色矩形边框",
        bool(candidates),
        f"符合位于页面四周边缘(边距≤{EDGE_TOL}cm)/蓝色单实线/0.5-3磅的外框 {len(candidates)} 个",
    )


def check_beaker_outline(ctx: EvaluationContext) -> RuleResult:
    # 细则区域：距左4.5–17cm，距上8–19.5cm
    REGION_L, REGION_T, REGION_R, REGION_B = 4.5, 8.0, 17.0, 19.5
    region_cx = (REGION_L + REGION_R) / 2
    region_h = REGION_B - REGION_T
    top_y_hi = REGION_T + region_h * 0.4     # 顶部椭圆允许出现的最深 y
    bottom_y_lo = REGION_T + region_h * 0.6  # 底部椭圆允许出现的最浅 y
    # 边线颜色：浅灰蓝色 —— 只要属于蓝色或灰色家族即视为合规
    def is_light_grey_blue(c: tuple[int, int, int] | None) -> bool:
        return is_blue(c) or is_grayish(c)
    # 填充：无填充、白色、或极浅蓝（透明效果）
    def is_no_or_light_fill(s: ShapeInfo):
        if not s.fill_colors:
            return True  # 无填充
        return all(
            is_whiteish(c) or (c[2] >= 200 and c[1] >= 200)
            for c in s.fill_colors
        )
    # 边线线宽在0.5–2磅内，非虚线单实线
    def has_outline_line(s: ShapeInfo):
        return (
            line_width_in_range(s.line_width, 0.5, 2.0, 0.0)
            and not s.dashed
            and is_light_grey_blue(s.line_color)
        )
    # 1. 上沿椭圆口：椭圆，中心在区域上部（占区域高 0–40%），宽>高（横向扁）
    top_ellipses = [
        s for s in ctx.ellipse_shapes
        if center_in(s, REGION_L, REGION_T, REGION_R, top_y_hi, tol=0.35)
        and s.w_abs > s.h_abs * 1.5
        and s.w_abs >= 4.0
        and has_outline_line(s)
    ]
    # 2. 斜壁筛选：位于区域内、颜色/线宽合规、且为"斜向"（既非纯水平也非纯垂直）
    #    - length >= 2cm
    #    - w_abs >= 0.15cm 且 h_abs >= 0.8cm 排除水平/极短竖线
    #    - 斜率 |h/w| ∈ [0.5, 15]（≈30°–86°），排除近水平/近垂直
    def is_slanted_wall(s: ShapeInfo) -> bool:
        if not (s.is_line or s.is_rect_like):
            return False
        if s.length < 2.0:
            return False
        if s.w_abs < 0.15 or s.h_abs < 0.8:
            return False
        slope = s.h_abs / max(s.w_abs, 1e-6)
        if not (0.5 <= slope <= 15.0):
            return False
        return has_outline_line(s)
    slanted_walls = [
        s for s in ctx.shapes
        if intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.2)
        and is_slanted_wall(s)
    ]
    # 3. 拆分左/右杯壁：以区域水平中线为界
    left_walls = [s for s in slanted_walls if s.cx < region_cx]
    right_walls = [s for s in slanted_walls if s.cx >= region_cx]
    # 4. 底部弧形：椭圆，中心在区域下部（占区域高 60%–100%），宽>高
    bottom_ellipses = [
        s for s in ctx.ellipse_shapes
        if center_in(s, REGION_L, bottom_y_lo, REGION_R, REGION_B, tol=0.4)
        and s.w_abs > s.h_abs
        and s.w_abs >= 3.0
        and has_outline_line(s)
    ]
    # 5. 上宽下窄验证：
    #    (a) 主路径 — 用左右壁 bbox 端点近似端口开口宽度。
    #        物理约束：左壁形似 "\"（上外下内），右壁形似 "/"（上外下内）,
    #        因此上开口宽 ≈ right_wall.x2 - left_wall.x1，
    #             下开口宽 ≈ right_wall.x1 - left_wall.x2。
    #        取每侧 length 最大的一条为代表；要求 top_gap ≥ bottom_gap + 0.3cm。
    #    (b) 退化路径 — 若主路径未成立(缺壁/形态异常), 用上/下椭圆宽度比较,
    #        top_ellipse.w_abs ≥ bottom_ellipse.w_abs + 0.3cm 亦视为通过。
    top_wider = False
    top_wider_detail = ""
    if left_walls and right_walls:
        lw = max(left_walls, key=lambda s: s.length)
        rw = max(right_walls, key=lambda s: s.length)
        top_gap = rw.x2 - lw.x1
        bottom_gap = rw.x1 - lw.x2
        if top_gap > 0 and bottom_gap > 0 and top_gap >= bottom_gap + 0.3:
            top_wider = True
            top_wider_detail = f"壁端口 上{top_gap:.2f}cm>下{bottom_gap:.2f}cm"
    if not top_wider and top_ellipses and bottom_ellipses:
        top_w = max(e.w_abs for e in top_ellipses)
        bot_w = max(e.w_abs for e in bottom_ellipses)
        if top_w >= bot_w + 0.3:
            top_wider = True
            top_wider_detail = f"椭圆 上{top_w:.2f}cm>下{bot_w:.2f}cm"
    if not top_wider and not top_wider_detail:
        top_wider_detail = "上宽下窄未验证"
    # 6. 整体填充检查：杯身主体形状（大面积）应无填充或极浅蓝透明
    body_shapes = [
        s for s in ctx.shapes
        if intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.2)
        and s.is_rect_like
        and s.w_abs >= 4.0 and s.h_abs >= 3.0
    ]
    fill_ok = not body_shapes or any(is_no_or_light_fill(s) for s in body_shapes)
    hit = (
        bool(top_ellipses)
        and bool(left_walls) and bool(right_walls)
        and bool(bottom_ellipses)
        and fill_ok
        and top_wider
    )
    return result(
        "第1页左侧烧杯外轮廓",
        hit,
        (
            f"上沿椭圆 {len(top_ellipses)}，左斜壁 {len(left_walls)}，"
            f"右斜壁 {len(right_walls)}，底部弧形 {len(bottom_ellipses)}，"
            f"填充合规={'是' if fill_ok else '否'}，"
            f"上宽下窄={'是' if top_wider else '否'}({top_wider_detail})"
        ),
    )


def check_beaker_top_ellipse(ctx: EvaluationContext) -> RuleResult:
    # 细则区域：距左3–18cm，距上7.2–13cm
    REGION_L, REGION_T, REGION_R, REGION_B = 3.0, 7.2, 18.0, 13.0
    # 边线颜色：浅蓝色 —— 蓝色家族即视为合规
    def is_light_blue(c: tuple[int, int, int] | None) -> bool:
        return is_blue(c)
    vals = []
    for s in ctx.ellipse_shapes:
        # 中心在指定区域内
        if not center_in(s, REGION_L, REGION_T, REGION_R, REGION_B, tol=0.35):
            continue
        # 横向扁椭圆：宽明显大于高
        if s.w_abs <= s.h_abs:
            continue
        # 边线为浅蓝色单实线
        if not is_light_blue(s.line_color):
            continue
        if s.dashed:
            continue
        # 线宽0.5磅–2磅（未显式设置按默认1.0磅处理）
        if not line_width_in_range(s.line_width, 0.5, 2.0, 0.0):
            continue
        vals.append(s)
    return result("第1页左侧烧杯顶部开口椭圆", bool(vals), f"符合横向扁椭圆 {len(vals)} 个")


def check_liquid_and_bubbles(ctx: EvaluationContext) -> RuleResult:
    # 细则区域：距左4–14.5cm，距上10–19cm
    REGION_L, REGION_T, REGION_R, REGION_B = 4.0, 10.0, 14.5, 19.0

    # 液体填充色：浅蓝色或淡青蓝色 —— 只要属于蓝色家族即视为合规
    def is_light_cyan_blue_fill(c: tuple[int, int, int] | None) -> bool:
        return is_blue(c) or is_cyan_or_whiteblue(c)

    # 1. 液体填充：在区域内，填充为浅蓝/淡青蓝，半透明（has_transparency），面积>=0.5
    fills = [
        s for s in ctx.shapes
        if intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.2)
        and s.has_transparency
        and any(is_light_cyan_blue_fill(c) for c in s.fill_colors)
        and s.box.area >= 0.5
    ]

    # 液面上沿边线颜色：浅蓝色 —— 只要属于蓝色家族即视为合规
    def is_surface_blue(c: tuple[int, int, int] | None) -> bool:
        return is_blue(c)

    # 2. 液面上沿：椭圆弧线，中心在区域上部(REGION_T 上下 ~2cm 范围)，横向扁，边线浅蓝色
    surface = [
        s for s in ctx.ellipse_shapes
        if center_in(s, REGION_L - 0.5, REGION_T - 2.0, REGION_R + 0.5, REGION_T + 1.5, tol=0.4)
        and s.w_abs > s.h_abs
        and s.w_abs >= 4.0
        and is_surface_blue(s.line_color)
    ]

    # 3. 气泡：在液体区域内（偏中下部），小椭圆/圆，填充为浅白或浅蓝白，直径<=0.5cm
    #    按 rubric 原文 "浅白色 或 浅蓝白色", 收窄判定:
    #      (a) 浅白  = 高亮度近白 (is_whiteish: min(RGB) ≥ 230);
    #      (b) 浅蓝白 = 高亮度浅蓝调 (三通道 ≥ 200 且 B 略高于 R, 允许极轻微蓝调).
    #    显式排除任意灰色 (原来 is_grayish 会把中/深灰误判为气泡, 已删除).
    def is_high_brightness_bluewhite(c: tuple[int, int, int] | None) -> bool:
        if c is None:
            return False
        r, g, b = c
        return min(r, g, b) >= 200 and b >= r and b >= g - 10

    def is_white_or_light_bluewhite(c: tuple[int, int, int] | None) -> bool:
        return is_whiteish(c) or is_high_brightness_bluewhite(c)

    bubbles = [
        s for s in ctx.ellipse_shapes
        if center_in(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.2)
        and s.w_abs <= 0.5 and s.h_abs <= 0.5
        and (
            not s.fill_colors
            or any(is_white_or_light_bluewhite(c) for c in s.fill_colors)
        )
    ]

    hit = bool(fills) and bool(surface) and len(bubbles) >= 3
    return result(
        "第1页左侧烧杯内部液面与液体填充",
        hit,
        f"浅蓝半透明液体形状 {len(fills)}，液面椭圆弧 {len(surface)}，气泡圆点 {len(bubbles)} 个",
    )


def check_sample_text(ctx: EvaluationContext) -> RuleResult:
    # 细则区域：距左3–9cm，距上9–13.5cm
    REGION_L, REGION_T, REGION_R, REGION_B = 3.0, 9.0, 9.0, 13.5
    # 1. 文本内容为"样品溶液"，中心在指定区域内
    vals = find_text_in(ctx, "样品溶液", REGION_L, REGION_T, REGION_R, REGION_B, tol=0.15)
    matched = []
    for s in vals:
        # 2. 字体：黑体、微软雅黑或相近中文字体（无衬线中文字体）
        font_ok = True  # 细则允许"相近中文字体"，有font_name时才做限制
        if s.font_name:
            fn_lower = s.font_name.lower()
            font_ok = any(k in fn_lower for k in (
                "黑体", "微软雅黑", "simhei", "microsoft yahei",
                "等线", "dengxian", "思源黑体", "noto sans sc",
                "heiti", "pingfang",
            ))
        # 3. 字号：14–16磅（严格，无额外容差）
        size_ok = s.font_size is not None and in_range(s.font_size, 14.0, 16.0, 0.0)
        # 4. 颜色：黑色或深灰色。
        #    要求 text_color 明确可解析(非 None) 且属于黑/深灰家族；
        #    未显式设置(a:srgbClr 缺失)的文本不再默认放行。
        #    build_context 已通过 _load_theme_colors 将 a:schemeClr(如 tx1/dk1) 反查
        #    theme1.xml 的 clrScheme 得到 RGB (纯 XML 解析, 无需 COM/LibreOffice)，
        #    因此使用主题色的文本也可正确解析为具体 RGB 再判定。
        color_ok = s.text_color is not None and (
            is_blackish(s.text_color) or is_dark(s.text_color)
        )
        # 5. 水平居中要求已从 rubric 中删除, 不再校验 text_align
        if font_ok and size_ok and color_ok:
            matched.append(s)
    hit = bool(matched)
    return result(
        "第1页左侧烧杯底部文本",
        hit,
        f"区域内样品溶液文本 {len(vals)} 个，全条件符合 {len(matched)} 个",
    )


def check_reference_electrode(ctx: EvaluationContext) -> RuleResult:
    # 细则区域：距左3.5–8.5cm，距上6–16.5cm
    REGION_L, REGION_T, REGION_R, REGION_B = 3.5, 6.0, 8.5, 16.5

    # 边线/填充颜色：绿色 —— 只要属于绿色家族即视为合规
    def is_green_shade(c: tuple[int, int, int] | None) -> bool:
        return is_green(c)

    # 向左倾斜判定:
    #   顶部在右、底部在左 → 电极从上到下向左倾斜。
    #   PowerPoint 中 xfrm@rot 正值=顺时针(屏幕 Y 向下), 因此对一个"本来竖直"的形状,
    #   rot<0 (逆时针) 表示顶端偏右、底端偏左, 即"向左倾斜"。反之 rot>0 表示向右倾斜。
    #   flipH 会反转左右, flipV 会反转上下, 二者叠加旋转的实际视觉倾斜方向。
    #   容差: 允许 -0.5° 以上(数值噪声)不算左倾, 上限-45°以内认为是"细长电极的合理左倾"。
    def tilts_left(s: ShapeInfo) -> tuple[bool, str]:
        r = s.rotation
        # flipH 反转水平, 相当于旋转方向翻转; flipV 反转垂直, 对左右倾斜方向也翻转
        eff = r
        if s.flip_h:
            eff = -eff
        if s.flip_v:
            eff = -eff
        # 归一化到 (-180, 180]
        eff = ((eff + 180.0) % 360.0) - 180.0
        # 对近 180° 的情况(整体上下颠倒但角度接近 180 而非 0), 折算等价倾角
        if eff > 90.0:
            eff -= 180.0
        elif eff < -90.0:
            eff += 180.0
        # eff∈(-90,90); <0 表示顶端偏右、底端偏左, 视觉上"向左倾斜"
        ok = -45.0 <= eff <= -0.5
        return ok, f"rot={r:.1f}°(等效{eff:.1f}°){',flipH' if s.flip_h else ''}{',flipV' if s.flip_v else ''}"

    # 主体：细长矩形，向左倾斜的近竖向, 在区域内, 绿色填充+绿色单实线, 线宽0.5–2磅
    body_shapes: list[ShapeInfo] = []
    body_debug: list[str] = []
    for s in ctx.rect_shapes:
        if not intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.3):
            continue
        if s.h_abs < 1.5:                       # 细长: 高度明显
            continue
        if s.h_abs <= s.w_abs * 2:              # 竖向为主
            continue
        if s.w_abs > 1.0:                        # 宽度细
            continue
        if s.dashed:
            continue
        if not is_green_shade(s.line_color):
            continue
        if not line_width_in_range(s.line_width, 0.5, 2.0, 0.0):
            continue
        if not any(is_green_shade(c) for c in s.fill_colors):
            continue
        left_ok, tag = tilts_left(s)
        if not left_ok:
            body_debug.append(f"#{s.index}非左倾({tag})")
            continue
        body_debug.append(f"#{s.index}左倾({tag})")
        body_shapes.append(s)

    # 顶部绿色连接头：在区域上部（y < 9cm），矩形或小方形，绿色填充+绿色边线
    top_connector = [
        s for s in ctx.rect_shapes
        if intersects_region(s, REGION_L, REGION_T, REGION_R, 9.0, 0.3)
        and s.box.area >= 0.05
        and s.box.area <= 0.8
        and is_green_shade(s.line_color)
        and any(is_green_shade(c) for c in s.fill_colors)
    ]

    # 底部伸入液体：主体下端（y > 10cm）有绿色对象（矩形或线段）
    bottom_part = [
        s for s in ctx.shapes
        if (s.is_rect_like or s.is_line)
        and intersects_region(s, REGION_L, 10.0, REGION_R, REGION_B, 0.3)
        and (is_green_shade(s.line_color) or any(is_green_shade(c) for c in s.fill_colors))
    ]

    # 如果没有旋转的矩形符合左倾, 尝试用"线段端点"兜底: 找区域内的绿色线, 顶端 x > 底端 x
    tilted_line_fallback: list[ShapeInfo] = []
    if not body_shapes:
        for s in ctx.line_shapes:
            if not intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.3):
                continue
            if not is_green_shade(s.line_color):
                continue
            if s.h_abs < 1.0:                   # 至少 1cm 的竖向跨度
                continue
            # 端点: 依据 flipV 决定顶端在 y1 还是 y2 (line 的 bbox 覆盖整根线,
            # 无 flipV 时约定顶端-左下, 但仅凭 bbox 无法确知斜率符号;
            # 用 flipH 作为斜率反转标识: OOXML 中 line 默认从左上到右下,
            # flipH=true 时变为从右上到左下, 即"向左倾斜")
            if s.h_abs <= s.w_abs:              # 需要竖向为主
                continue
            left_lean = s.flip_h != s.flip_v    # 顶端在右下端在左
            if left_lean:
                tilted_line_fallback.append(s)

    body_ok = bool(body_shapes) or bool(tilted_line_fallback)
    hit = body_ok and bool(top_connector) and bool(bottom_part)
    return result(
        "第1页左侧绿色参比电极",
        hit,
        (
            f"左倾主体 {len(body_shapes)} 个"
            f"{'(+线段兜底 ' + str(len(tilted_line_fallback)) + ')' if tilted_line_fallback else ''}, "
            f"顶部连接头 {len(top_connector)}, 底部入液 {len(bottom_part)}"
            f"{'; 未采纳: ' + '; '.join(body_debug[:3]) if not body_shapes and body_debug else ''}"
        ),
    )


def check_working_electrode(ctx: EvaluationContext) -> RuleResult:
    # 细则区域：距左6–12cm，距上3.6–17.5cm
    REGION_L, REGION_T, REGION_R, REGION_B = 6.0, 3.6, 12.0, 17.5

    def is_dark_blue_black(c: tuple[int, int, int] | None) -> bool:
        """深蓝黑色 —— 只要属于深色（黑/深灰/深蓝）即视为合规"""
        return is_blackish(c) or is_dark(c) or is_blue(c) and is_dark(c)

    def is_blackish_grey(c: tuple[int, int, int] | None) -> bool:
        """黑色或深灰色 —— 只要属于黑色/深色/灰色家族即视为合规"""
        return is_blackish(c) or is_dark(c) or is_grayish(c)

    def is_white_rod(c: tuple[int, int, int] | None) -> bool:
        """白色杆体"""
        return is_whiteish(c)

    def is_light_blue_sleeve(c: tuple[int, int, int] | None) -> bool:
        """浅蓝色透明套筒 —— 只要属于蓝色家族即视为合规"""
        return is_blue(c) or is_cyan_or_whiteblue(c)

    all_in = [s for s in ctx.shapes if intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.3)]
    rects_in = [s for s in all_in if s.is_rect_like]

    # 1. 近垂直整体：细长竖向矩形，h>=4cm，h>w*3
    vertical_body = [
        s for s in rects_in
        if s.h_abs >= 4.0 and s.h_abs > s.w_abs * 3
    ]

    # 2. 顶部深蓝黑色圆柱头：区域上部（y<6.5cm），深色填充矩形
    top_head = [
        s for s in rects_in
        if s.cy <= 6.5
        and any(is_dark_blue_black(c) for c in s.fill_colors)
    ]

    # 3. 黑色固定环：深色填充小矩形，在中部（y 6–8cm），宽度比主体略宽
    dark_ring = [
        s for s in rects_in
        if 6.0 <= s.cy <= 8.5
        and any(is_blackish_grey(c) for c in s.fill_colors)
        and s.w_abs >= 0.3
    ]

    # 4. 白色杆体
    white_rod = [
        s for s in rects_in
        if s.h_abs >= 2.0
        and any(is_white_rod(c) for c in s.fill_colors)
    ]

    # 5. 浅蓝色套筒（透明效果）：浅蓝填充，在中下部
    blue_sleeve = [
        s for s in rects_in
        if (s.has_transparency or any(is_light_blue_sleeve(c) for c in s.fill_colors))
        and any(is_light_blue_sleeve(c) for c in s.fill_colors)
        and s.h_abs >= 0.5
    ]

    # 6. 边线深灰或黑色单实线，线宽0.5–2磅
    has_valid_outline = any(
        s for s in rects_in
        if not s.dashed
        and is_blackish_grey(s.line_color)
        and line_width_in_range(s.line_width, 0.5, 2.0, 0.0)
    )

    # 7. 下端伸入液体: rubric 明确要求"下端伸入液体中"。液面在 y≈9cm(参考液体检查区域上沿),
    #    要求主体或套筒的底端 y2 至少达到 9cm, 即末端进入液体区。
    #    优先看主体(vertical_body)/套筒(blue_sleeve); 若都缺, 兜底看整个 rects_in 内最深 y2。
    LIQUID_Y = 9.0
    dip_candidates = vertical_body + blue_sleeve + white_rod
    if dip_candidates:
        max_bottom = max(s.y2 for s in dip_candidates)
    elif rects_in:
        max_bottom = max(s.y2 for s in rects_in)
    else:
        max_bottom = -1.0
    dips_into_liquid = max_bottom >= LIQUID_Y - 0.15  # 0.15cm 容差

    hit = (
        bool(vertical_body) and bool(top_head) and bool(dark_ring)
        and bool(white_rod) and bool(blue_sleeve)
        and has_valid_outline and dips_into_liquid
    )
    return result(
        "第1页左侧中间白蓝工作电极",
        hit,
        (
            f"竖向主体 {len(vertical_body)}，深蓝黑圆柱头 {len(top_head)}，黑色固定环 {len(dark_ring)}，"
            f"白色杆 {len(white_rod)}，浅蓝套筒 {len(blue_sleeve)}，"
            f"边线合规={'是' if has_valid_outline else '否'}，"
            f"下端入液(y2max={max_bottom:.2f}cm≥{LIQUID_Y}cm)={'是' if dips_into_liquid else '否'}"
        ),
    )



def check_counter_electrode(ctx: EvaluationContext) -> RuleResult:
    # 细则区域：距左8.5–13.5cm，距上6.5–16.5cm
    REGION_L, REGION_T, REGION_R, REGION_B = 8.5, 6.5, 13.5, 16.5

    # 边线/填充颜色：红色到深红色 —— 只要属于红色家族即视为合规
    def is_red_shade(c: tuple[int, int, int] | None) -> bool:
        return is_red(c)

    # 向右倾斜判定:
    #   顶部在左、底部在右 → 电极从上到下向右倾斜。
    #   PowerPoint xfrm@rot 正值=顺时针(屏幕 Y 向下), 对本来竖直的形状:
    #     rot>0 (顺时针) 表示顶端偏左、底端偏右 → 视觉"向右倾斜"。
    #   flipH / flipV 各自等价于旋转方向取反。
    #   容差: eff∈(-90,90), 认为 +0.5° ≤ eff ≤ +45° 是合规右倾。
    def tilts_right(s: ShapeInfo) -> tuple[bool, str]:
        r = s.rotation
        eff = r
        if s.flip_h:
            eff = -eff
        if s.flip_v:
            eff = -eff
        eff = ((eff + 180.0) % 360.0) - 180.0
        if eff > 90.0:
            eff -= 180.0
        elif eff < -90.0:
            eff += 180.0
        ok = 0.5 <= eff <= 45.0
        return ok, f"rot={r:.1f}°(等效{eff:.1f}°){',flipH' if s.flip_h else ''}{',flipV' if s.flip_v else ''}"

    # 主体：细长矩形，向右倾斜的近竖向, 在区域内, 红色填充+红色单实线, 线宽0.5–2磅
    body_shapes: list[ShapeInfo] = []
    body_debug: list[str] = []
    for s in ctx.rect_shapes:
        if not intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.3):
            continue
        if s.h_abs < 1.5:
            continue
        if s.h_abs <= s.w_abs * 2:
            continue
        if s.w_abs > 1.0:
            continue
        if s.dashed:
            continue
        if not is_red_shade(s.line_color):
            continue
        if not line_width_in_range(s.line_width, 0.5, 2.0, 0.0):
            continue
        if not any(is_red_shade(c) for c in s.fill_colors):
            continue
        right_ok, tag = tilts_right(s)
        if not right_ok:
            body_debug.append(f"#{s.index}非右倾({tag})")
            continue
        body_debug.append(f"#{s.index}右倾({tag})")
        body_shapes.append(s)

    # 顶部红色连接头：在区域上部（y < 8.5cm），矩形或小方形，红色填充+红色边线，线宽0.5–2磅
    top_connector = [
        s for s in ctx.rect_shapes
        if intersects_region(s, REGION_L, REGION_T, REGION_R, 8.5, 0.3)
        and 0.05 <= s.box.area <= 0.8
        and not s.dashed
        and is_red_shade(s.line_color)
        and line_width_in_range(s.line_width, 0.5, 2.0, 0.0)
        and any(is_red_shade(c) for c in s.fill_colors)
    ]

    # 底部伸入液体：区域下部（y > 10cm）有红色矩形或线段
    bottom_part = [
        s for s in ctx.shapes
        if (s.is_rect_like or s.is_line)
        and intersects_region(s, REGION_L, 10.0, REGION_R, REGION_B, 0.3)
        and (is_red_shade(s.line_color) or any(is_red_shade(c) for c in s.fill_colors))
    ]

    # 线段兜底: 若无右倾矩形, 找区域内红色 line, 依据 flipH/flipV 判断斜率符号
    tilted_line_fallback: list[ShapeInfo] = []
    if not body_shapes:
        for s in ctx.line_shapes:
            if not intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.3):
                continue
            if not is_red_shade(s.line_color):
                continue
            if s.h_abs < 1.0:
                continue
            if s.h_abs <= s.w_abs:
                continue
            # OOXML line 默认左上→右下(顶端在左, 底端在右) = "向右倾斜";
            # flipH XOR flipV 为 True 时反转为左倾, 因此右倾即 flip_h == flip_v。
            right_lean = s.flip_h == s.flip_v
            if right_lean:
                tilted_line_fallback.append(s)

    body_ok = bool(body_shapes) or bool(tilted_line_fallback)
    hit = body_ok and bool(top_connector) and bool(bottom_part)
    return result(
        "第1页右侧红色对电极",
        hit,
        (
            f"右倾主体 {len(body_shapes)} 个"
            f"{'(+线段兜底 ' + str(len(tilted_line_fallback)) + ')' if tilted_line_fallback else ''}, "
            f"顶部连接头 {len(top_connector)}, 底部入液 {len(bottom_part)}"
            f"{'; 未采纳: ' + '; '.join(body_debug[:3]) if not body_shapes and body_debug else ''}"
        ),
    )



def check_reference_label_arrow(ctx: EvaluationContext) -> RuleResult:
    """第1页左侧"参比电极"文本 + 右侧黑色水平单箭头。
    rubric:
      文本区域: 距左 2–4.4cm, 距上 9.2–10.4cm; 内容 "参比电极"(兼容"参考电极");
      字号:    13–15 磅(严格, 无额外容差);
      字体:    黑体/微软雅黑/相近中文字体;
      颜色:    黑色或深灰色 —— text_color 必须明确可解析。
      箭头区域: 距左 4.2–6.2cm, 距上 9.4–9.8cm;
      箭头:    黑色水平单箭头, 长度 1.4–1.6cm, 线宽 1–1.5 磅, 箭头朝右, 末端指向绿色参比电极。
    """
    # 文本区域：距左2–4.4cm, 距上9.2–10.4cm; 文本为"参比电极"(兼容"参考电极")
    texts: list[ShapeInfo] = find_text_in(ctx, ["参比电极", "参考电极"], 2.0, 9.2, 4.4, 10.4, tol=0.15)
    valid_texts: list[ShapeInfo] = []
    for s in texts:
        # 字体: 黑体/微软雅黑或相近中文字体 —— 无 font_name 时允许(细则允许"相近中文字体")
        font_ok = True
        if s.font_name:
            fn = s.font_name.lower()
            font_ok = any(k in fn for k in (
                "黑体", "微软雅黑", "simhei", "microsoft yahei",
                "等线", "dengxian", "思源黑体", "noto sans sc",
                "heiti", "pingfang",
            ))
        # 字号: 13–15 磅
        size_ok = s.font_size is not None and in_range(s.font_size, 13.0, 15.0, 0.0)
        # 颜色: 明确可解析 且 黑色/深灰
        color_ok = s.text_color is not None and (
            is_blackish(s.text_color) or is_dark(s.text_color)
        )
        if font_ok and size_ok and color_ok:
            valid_texts.append(s)

    # 箭头：区域 4.2–6.2cm × 9.4–9.8cm, 黑色水平单箭头, 长度 1.4–1.6cm, 线宽 1–1.5 磅, 朝右
    def _is_black_line(c: Optional[tuple[int, int, int]]) -> bool:
        return is_blackish(c) or is_dark(c)

    arrows = line_in_region(
        ctx, 4.2, 9.4, 6.2, 9.8,
        color_pred=_is_black_line,
        width_range=(1.0, 1.5), width_tol=0.0,
        dashed=False, arrow=True, horizontal=True, direction="right",
        min_len=1.4, max_len=1.6,
    )

    hit = bool(valid_texts) and bool(arrows)
    return result(
        '第1页左侧"参比电极"文本',
        hit,
        f"区域内文本 {len(texts)} 个，全条件符合 {len(valid_texts)} 个，右向箭头 {len(arrows)} 条",
    )


def check_working_label(ctx: EvaluationContext) -> RuleResult:
    # 细则区域：距左6.4–13cm，距上9–14cm
    REGION_L, REGION_T, REGION_R, REGION_B = 6.4, 9.0, 13.0, 14.0
    texts: list[ShapeInfo] = find_text_in(ctx, "工作电极", REGION_L, REGION_T, REGION_R, REGION_B, tol=0.15)
    matched: list[ShapeInfo] = []
    for s in texts:
        # 字体：黑体、微软雅黑或相近中文字体
        font_ok = True
        if s.font_name:
            fn = s.font_name.lower()
            font_ok = any(k in fn for k in (
                "黑体", "微软雅黑", "simhei", "microsoft yahei",
                "等线", "dengxian", "思源黑体", "noto sans sc",
                "heiti", "pingfang",
            ))
        # 字号14–22磅（严格无容差）
        size_ok = s.font_size is not None and in_range(s.font_size, 14.0, 22.0, 0.0)
        # 颜色黑色或深灰色 —— 明确可解析(theme 已在 build_context 反查); 无颜色不再默认放行
        color_ok = s.text_color is not None and (
            is_blackish(s.text_color) or is_dark(s.text_color)
        )
        # rubric 已删除"水平居中", 不再校验 text_align
        if font_ok and size_ok and color_ok:
            matched.append(s)
    hit = bool(matched)
    return result('第1页中部"工作电极"文本', hit, f"区域内文本 {len(texts)} 个，全条件符合 {len(matched)} 个")


def check_counter_label_arrow(ctx: EvaluationContext) -> RuleResult:
    """第1页右侧"对电极"文本 + 左侧黑色水平单箭头。
    rubric:
      文本区域: 距左 10.5–12.7cm, 距上 9.1–10.3cm; 内容 "对电极";
      字号:    14–16 磅(严格, 无额外容差);
      字体:    黑体/微软雅黑/相近中文字体;
      颜色:    黑色或深灰色 —— text_color 必须明确可解析。
      箭头区域: 距左 9.2–11.2cm, 距上 9.4–9.9cm;
      箭头:    黑色水平单箭头, 长度 1.4–1.6cm, 线宽 1–1.5 磅, 箭头朝右, 末端指向红色对电极。
    """
    # 文本区域：距左10.5–12.7cm, 距上9.1–10.3cm; 文本为"对电极"
    TEXT_L, TEXT_T, TEXT_R, TEXT_B = 10.5, 9.1, 12.7, 10.3
    texts: list[ShapeInfo] = find_text_in(ctx, "对电极", TEXT_L, TEXT_T, TEXT_R, TEXT_B, tol=0.15)
    valid_texts: list[ShapeInfo] = []
    for s in texts:
        # 字体：黑体、微软雅黑或相近中文字体
        font_ok = True
        if s.font_name:
            fn = s.font_name.lower()
            font_ok = any(k in fn for k in (
                "黑体", "微软雅黑", "simhei", "microsoft yahei",
                "等线", "dengxian", "思源黑体", "noto sans sc",
                "heiti", "pingfang",
            ))
        # 字号14–16 磅（严格无容差）
        size_ok = s.font_size is not None and in_range(s.font_size, 14.0, 16.0, 0.0)
        # 颜色: 明确可解析 且 黑/深灰
        color_ok = s.text_color is not None and (
            is_blackish(s.text_color) or is_dark(s.text_color)
        )
        if font_ok and size_ok and color_ok:
            valid_texts.append(s)

    # 箭头：区域 9.2–11.2cm × 9.4–9.9cm, 黑色水平单箭头, 长度 1.4–1.6cm, 线宽 1–1.5 磅, 朝右
    def _is_black_line(c: Optional[tuple[int, int, int]]) -> bool:
        return is_blackish(c) or is_dark(c)

    arrows = line_in_region(
        ctx, 9.2, 9.4, 11.2, 9.9,
        color_pred=_is_black_line,
        width_range=(1.0, 1.5), width_tol=0.0,
        dashed=False, arrow=True, horizontal=True, direction="right",
        min_len=1.4, max_len=1.6,
    )

    hit = bool(valid_texts) and bool(arrows)
    return result(
        '第1页右侧"对电极"文本',
        hit,
        f"区域内文本 {len(texts)} 个，全条件符合 {len(valid_texts)} 个，右向箭头 {len(arrows)} 条",
    )



def check_analyte_input(ctx: EvaluationContext) -> RuleResult:
    """第1页顶部"分析物输入"文本 + 下方蓝色竖向虚线箭头。
    rubric:
      文本区域: 距左 7–10.3cm, 距上 1.0–1.8cm; 内容 "分析物输入";
      字号:    14–16 磅(严格, 无额外容差);
      字体:    黑体/微软雅黑/相近中文字体;
      颜色:    蓝色 —— text_color 必须明确可解析。
      箭头区域: 距左 7.2–8.4cm, 距上 1.5–3.8cm;
      箭头:    蓝色竖向虚线单箭头, 线宽 1–1.5 磅, 箭头朝下, 末端指向工作电极顶部。
    箭头支持两条判定路径 (rubric 只要求视觉上的"虚线单箭头", 具体绘制方式不限):
      1) 原生箭头: line/连接符设置 headEnd/tailEnd 并设为虚线, 由 line_in_region 校验;
      2) 分段虚线: 多段蓝色短竖线 + 下方三角/downArrow 组合(与 check_signal_box 同构),
         分段短竖线线宽 1.0–1.5 磅, 三角/downArrow 必须"朝下"且位于短竖线组下方。
    """
    # 文本区域：距左7–10.3cm, 距上1.0–1.8cm; 文本为"分析物输入"
    TEXT_L, TEXT_T, TEXT_R, TEXT_B = 7.0, 1.0, 10.3, 1.8
    texts: list[ShapeInfo] = find_text_in(ctx, "分析物输入", TEXT_L, TEXT_T, TEXT_R, TEXT_B, tol=0.15)
    valid_texts: list[ShapeInfo] = []
    for s in texts:
        # 字体：黑体、微软雅黑或相近中文字体
        font_ok = True
        if s.font_name:
            fn = s.font_name.lower()
            font_ok = any(k in fn for k in (
                "黑体", "微软雅黑", "simhei", "microsoft yahei",
                "等线", "dengxian", "思源黑体", "noto sans sc",
                "heiti", "pingfang",
            ))
        # 字号14–16 磅（严格无容差）
        size_ok = s.font_size is not None and in_range(s.font_size, 14.0, 16.0, 0.0)
        # 颜色为蓝色 —— 明确可解析(schemeClr 已在 first_xml_rgb 里反查为 RGB)
        color_ok = s.text_color is not None and is_blue(s.text_color)
        if font_ok and size_ok and color_ok:
            valid_texts.append(s)

    # 箭头区域：7.2–8.4cm × 1.5–3.8cm
    ARROW_L, ARROW_T, ARROW_R, ARROW_B = 7.2, 1.5, 8.4, 3.8

    # 路径 1: 原生蓝色竖向虚线单箭头, 线宽 1–1.5 磅, 朝下
    native_arrows = line_in_region(
        ctx, ARROW_L, ARROW_T, ARROW_R, ARROW_B,
        color_pred=is_blue,
        width_range=(1.0, 1.5), width_tol=0.0,
        dashed=True, arrow=True, vertical=True, direction="down",
    )

    # 路径 2: 分段虚线 —— 多段蓝色短竖线 + 下方三角/downArrow 组合
    # 短竖线线宽严格 1.0–1.5 磅(与 rubric 一致, 不放宽到 1.8)
    blue_vsegs = [
        s for s in ctx.line_shapes
        if intersects_region(s, ARROW_L, ARROW_T, ARROW_R, ARROW_B, 0.05)
        and is_blue(s.line_color)
        and s.w_abs <= 0.05 and s.h_abs >= 0.05
        and line_width_in_range(s.line_width, 1.0, 1.5, 0.0)
    ]
    # 竖线组底端 y (无竖线时用 ARROW_T 兜底)
    vseg_bottom_y = max((s.y2 for s in blue_vsegs), default=ARROW_T)

    def tip_points_down(t: ShapeInfo) -> bool:
        """三角尖朝下判定:
          - downArrow 预设几何默认朝下, 允许 rot 偏离 ≤ 45°;
          - triangle 默认尖角朝上, 需 flipV=True 或 |rot| ≥ 135° 才朝下;
          - 其他几何一律视为不朝下。
        """
        r = getattr(t, "rotation", 0.0) or 0.0
        flip_v = getattr(t, "flip_v", False)
        eff = ((r + 180.0) % 360.0) - 180.0  # 归一化到 (-180, 180]
        if t.geom == "downArrow":
            return abs(eff) <= 45.0
        if t.geom == "triangle":
            if flip_v:
                return True
            return abs(eff) >= 135.0
        return False

    raw_tips = [
        s for s in ctx.shapes
        if s.geom in ("triangle", "downArrow")
        # 三角/箭头允许略越 ARROW_B (箭尖伸到工作电极顶部之上，工作电极顶 y≈3.6cm)
        and intersects_region(s, ARROW_L, ARROW_T, ARROW_R, ARROW_B + 0.4, 0.1)
        and (is_blue(s.line_color) or any(is_blue(c) for c in s.fill_colors))
    ]
    arrow_tips = [
        s for s in raw_tips
        if tip_points_down(s)
        # 箭尖 y (下沿) 位于短竖线组底端之下(容差 0.35cm)
        and s.y2 >= vseg_bottom_y - 0.35
    ]
    seg_arrow_ok = len(blue_vsegs) >= 2 and bool(arrow_tips)

    arrows_ok = bool(native_arrows) or seg_arrow_ok
    hit = bool(valid_texts) and arrows_ok
    return result(
        '第1页顶部"分析物输入"文本',
        hit,
        f"区域内文本 {len(texts)} 个，全条件符合 {len(valid_texts)} 个，"
        + f"原生向下箭头 {len(native_arrows)}，分段竖线 {len(blue_vsegs)}/朝下箭尖 {len(arrow_tips)}",
    )



def check_electrode_connection_lines(ctx: EvaluationContext) -> RuleResult:
    # 三条连接线均由多段短实线组成弧形曲线（Office中曲线离散化），颜色各异
    # 细则要求：各自区域内，对应颜色，单实线（非虚线），线宽0.5–3磅

    def collect_arc_segs(region, color_pred):
        """收集区域内指定颜色的弧形线段（单实线，线宽0.5–3磅）"""
        return [
            s for s in ctx.line_shapes
            if e_intersects(s, region)
            and color_pred(s.line_color)
            and not s.dashed
            and line_width_in_range(s.line_width, 0.5, 3.0, 0.0)
        ]

    def e_intersects(s, region):
        l, t, r, b = region
        return intersects_region(s, l, t, r, b, 0.1)

    # 绿色参比电极连接线：距左3.5–22cm，距上5.2–15cm，绿色弧形单实线
    green_region = (3.5, 5.2, 22.0, 15.0)
    green_segs = collect_arc_segs(green_region, is_green)
    # 起点在参比电极顶部附近（x∈[3.5,8.5], y∈[5.2,11]，与参比电极新区域一致）
    green_start = [s for s in green_segs if intersects_region(s, 3.5, 5.2, 8.5, 11.0, 0.2)]
    # 终点在仪器左上/左后侧（x∈[12,22], y∈[7,15]）
    green_end = [s for s in green_segs if intersects_region(s, 12.0, 7.0, 22.0, 15.0, 0.2)]
    green_ok = len(green_segs) >= 3 and bool(green_start) and bool(green_end)

    # 黑色工作电极连接线：距左7–22cm，距上3–14cm，黑色弧形单实线
    black_region = (7.0, 3.0, 22.0, 14.0)
    black_segs = collect_arc_segs(black_region, lambda c: is_blackish(c) or is_dark(c))
    # 起点在工作电极顶部附近（x∈[6,12], y∈[3,6]，与工作电极新区域顶部一致）
    black_start = [s for s in black_segs if intersects_region(s, 6.0, 3.0, 12.0, 6.0, 0.2)]
    # 终点在仪器左上/左后侧（x∈[12,22], y∈[7,14]）
    black_end = [s for s in black_segs if intersects_region(s, 12.0, 7.0, 22.0, 14.0, 0.2)]
    black_ok = len(black_segs) >= 3 and bool(black_start) and bool(black_end)

    # 红色对电极连接线：距左9.5–22cm，距上5.2–14cm，红色弧形单实线
    red_region = (9.5, 5.2, 22.0, 14.0)
    red_segs = collect_arc_segs(red_region, is_red)
    # 起点在对电极顶部附近（x∈[8.5,13.5], y∈[5.2,10]，与对电极新区域顶部一致）
    red_start = [s for s in red_segs if intersects_region(s, 8.5, 5.2, 13.5, 10.0, 0.2)]
    # 终点在仪器左上/左后侧（x∈[12,22], y∈[7,14]）
    red_end = [s for s in red_segs if intersects_region(s, 12.0, 7.0, 22.0, 14.0, 0.2)]
    red_ok = len(red_segs) >= 3 and bool(red_start) and bool(red_end)

    hit = green_ok and black_ok and red_ok
    return result(
        "第1页三条电极连接线",
        hit,
        f"绿色线段 {len(green_segs)}(起{len(green_start)}/终{len(green_end)})，"
        f"黑色线段 {len(black_segs)}(起{len(black_start)}/终{len(black_end)})，"
        f"红色线段 {len(red_segs)}(起{len(red_start)}/终{len(red_end)})",
    )


def check_analyzer_shell(ctx: EvaluationContext) -> RuleResult:
    """第1页中部电化学分析仪外壳。
    修订后 rubric 要素:
      区域: 距左 13.0–31cm, 距上 8.5–18cm;
      形态: 三个及以上的四边形及其他图形组合而成 或 整体为带透视效果的长方体仪器;
      前面: 蓝色;
      顶/侧面: 浅灰蓝到灰白色渐变 或 单纯浅灰蓝/灰白;
      底部: 深色底座阴影;
      边线: 深蓝色或灰蓝色单实线, 线宽 0.5–3 磅。
    透视路径必须同时确认 "顶面" 与 "侧面" 都存在(真实透视几何: parallelogram / trapezoid),
    不能靠普通矩形凑数; 若达不成, 走图形组合兜底(3 个及以上的合规四边形/椭圆等)。
    """
    # 细则区域：距左13.0–31cm，距上8.5–18cm
    REGION_L, REGION_T, REGION_R, REGION_B = 13.0, 8.5, 31.0, 18.0

    def is_grey_blue_or_grey_white(c: tuple[int, int, int] | None) -> bool:
        # 灰蓝/灰白 —— 只要属于蓝色/灰色/白色家族即视为合规
        return is_blue(c) or is_grayish(c) or is_whiteish(c)

    def is_deep_blue_or_grey_blue(c: tuple[int, int, int] | None) -> bool:
        # 深蓝/灰蓝 —— 只要属于蓝色或灰色家族即视为合规
        return is_blue(c) or is_grayish(c)

    def has_shell_outline(s: ShapeInfo) -> bool:
        # 线宽 0.5–3 磅, 单实线, 边线为深蓝/灰蓝
        return (
            not s.dashed
            and line_width_in_range(s.line_width, 0.5, 3.0, 0.0)
            and is_deep_blue_or_grey_blue(s.line_color)
        )

    region_shapes = [s for s in ctx.shapes if intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.15)]

    # 前面板：蓝色长方形/圆角矩形，位于仪器正面
    front_panel = [
        s for s in ctx.rect_shapes
        if intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.15)
        and s.w_abs >= 4.5
        and s.h_abs >= 2.0
        and any(is_blue(c) for c in s.fill_colors)
        and has_shell_outline(s)
    ]

    # 透视面 (真实透视几何: parallelogram / trapezoid) —— 浅灰蓝/灰白填充, 合规边线
    perspective_faces = [
        s for s in ctx.rect_shapes
        if intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.15)
        and s.geom in {"parallelogram", "trapezoid"}
        and any(is_grey_blue_or_grey_white(c) for c in s.fill_colors)
        and has_shell_outline(s)
    ]

    # 分别判定 "顶面" 与 "侧面":
    #   顶面: 位于前面板上方 —— cy 明显低于前面板 cy(y 坐标向下为正, 故 cy 更小),
    #         或底沿接近前面板顶沿 (y2 与前面板 y1 相近, ±1.2cm 容差)。
    #   侧面: 位于前面板左侧或右侧 —— cx 明显偏出前面板一侧,
    #         或其一竖直边接近前面板的左右边 (±1.2cm 容差)。
    #   若无 front_panel, 退化为按位置相对 REGION 中心分类。
    def is_top_face(face: ShapeInfo, panels: list[ShapeInfo]) -> bool:
        if panels:
            for fp in panels:
                # 底沿贴合前面板顶沿
                if abs(face.y2 - fp.y1) <= 1.2:
                    return True
                # 中心明显在前面板中心之上, 且横向有重叠
                if face.cy + 0.5 <= fp.cy and not (face.x2 < fp.x1 or face.x1 > fp.x2):
                    return True
            return False
        # 无前面板时: 面中心 y 在区域上半 (相对 REGION 中心线)
        return face.cy <= (REGION_T + REGION_B) / 2

    def is_side_face(face: ShapeInfo, panels: list[ShapeInfo]) -> bool:
        if panels:
            for fp in panels:
                # 竖直边贴合前面板左右边
                if abs(face.x1 - fp.x2) <= 1.2 or abs(face.x2 - fp.x1) <= 1.2:
                    return True
                # 中心明显在前面板中心之外, 且竖向有重叠
                horiz_off = (face.cx > fp.x2 + 0.5) or (face.cx < fp.x1 - 0.5)
                if horiz_off and not (face.y2 < fp.y1 or face.y1 > fp.y2):
                    return True
            return False
        # 无前面板: 面中心 x 偏出区域中心
        cx_mid = (REGION_L + REGION_R) / 2
        return abs(face.cx - cx_mid) >= 3.0

    top_faces = [f for f in perspective_faces if is_top_face(f, front_panel)]
    side_faces = [f for f in perspective_faces if is_side_face(f, front_panel)]

    # 底部深色底座阴影：位于仪器底部，深色，横向较宽
    shadow = [
        s for s in region_shapes
        if intersects_region(s, REGION_L, REGION_B - 3.5, REGION_R, REGION_B + 0.5, 0.15)
        and s.w_abs >= 3.0
        and s.h_abs <= 1.2
        and (is_dark(s.line_color) or any(is_dark(c) for c in s.fill_colors))
    ]

    # rubric 新增: "三个及以上的四边形及其他图形组合而成" —— 作为长方体/透视之外的兜底判据
    # 组合体判定: 边线合规的多边形/矩形/椭圆等在区域内累计 >= 3 个即视为"图形组合"
    combo_shapes = [
        s for s in region_shapes
        if (s.is_rect_like or s.is_ellipse)
        and has_shell_outline(s)
    ]
    combo_ok = len(combo_shapes) >= 3

    # 长方体透视判据: 蓝色前面板 + 同时具备真实"顶面"与"侧面"透视几何。
    #   不再接受普通矩形作 grey_body_faces 凑数。
    perspective_ok = bool(front_panel) and bool(top_faces) and bool(side_faces)

    # 满足两种形态之一 + 底部阴影 即视为通过
    hit = (perspective_ok or combo_ok) and bool(shadow)
    return result(
        "第1页中部电化学分析仪外壳",
        hit,
        (
            f"蓝色前面板 {len(front_panel)}，透视面 {len(perspective_faces)}"
            f"(顶面 {len(top_faces)}/侧面 {len(side_faces)})，"
            f"图形组合 {len(combo_shapes)}，底座阴影 {len(shadow)}"
        ),
    )


def check_analyzer_screen(ctx: EvaluationContext) -> RuleResult:
    """第1页中部电化学分析仪前屏幕。
    修订后 rubric:
      区域: 距左 16–28cm, 距上 10–18cm;
      形状: 蓝色矩形;
      边框: 白色或浅灰白色, 单实线, 线宽 0.5–3 磅;
      内部: 白色坐标轴+白色折线图 或 组合图形 或 图片。
    """
    # 细则区域：距左16–28cm，距上10–18cm
    REGION_L, REGION_T, REGION_R, REGION_B = 16.0, 10.0, 28.0, 18.0

    def is_white_or_light(c: tuple[int, int, int] | None) -> bool:
        return is_whiteish(c) or is_grayish(c)

    # 蓝色矩形显示屏：蓝色填充，边框为白色或浅灰白，线宽0.5–3磅，单实线
    screens = [
        s for s in ctx.rect_shapes
        if intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.15)
        and s.w_abs >= 1.5 and s.h_abs >= 1.0
        and any(is_blue(c) for c in s.fill_colors)
        and (s.line_color is None or is_white_or_light(s.line_color))
        and not s.dashed
        and line_width_in_range(s.line_width, 0.5, 3.0, 0.0)
    ]

    # 内部内容: 三种来源之一
    #   (a) 白色/浅色 "坐标轴 + 折线图" 组合:
    #        · 至少 1 条水平坐标轴 (h_abs≈0, w_abs≥0.4cm)
    #        · 至少 1 条竖直坐标轴 (w_abs≈0, h_abs≥0.4cm)
    #        · 至少 1 条折线/曲线段 (斜线: 非纯水平且非纯竖直, 或多段离散化短线共 ≥2 段)
    #   (b) 组合图形: 屏幕区域内除屏幕本身外, 有 >=2 个可见形状 (rect_like / ellipse / line);
    #   (c) 图片: 屏幕区域内有 has_blip=True 的形状。
    inner_L = REGION_L + 0.2
    inner_T = REGION_T + 0.2
    inner_R = REGION_R - 0.1
    inner_B = REGION_B - 0.2
    white_lines = line_in_region(
        ctx,
        inner_L, inner_T, inner_R, inner_B,
        color_pred=lambda c: is_whiteish(c) or is_cyan_or_whiteblue(c),
        width_range=(0.5, 3.0),
        dashed=False,
        min_len=0.1,
        width_tol=0.0,
        region_tol=0.05,
    )
    # 分类: 水平轴 / 竖直轴 / 折线段
    axis_thresh = 0.06  # cm, 视为"纯水平/纯竖直"的偏差容差
    h_axes = [s for s in white_lines if s.h_abs <= axis_thresh and s.w_abs >= 0.4]
    v_axes = [s for s in white_lines if s.w_abs <= axis_thresh and s.h_abs >= 0.4]
    diag_segs = [
        s for s in white_lines
        if s.w_abs > axis_thresh and s.h_abs > axis_thresh
    ]
    # 折线判据: 至少 1 条明确的斜线段, 或水平/竖直短碎段各有一些且总数 ≥ 2 段视为离散化折线
    non_axis_segs = [
        s for s in white_lines
        if s not in h_axes and s not in v_axes
    ]
    polyline_ok = bool(diag_segs) or len(non_axis_segs) >= 2
    lines_ok = bool(h_axes) and bool(v_axes) and polyline_ok

    screen_ids = {id(s) for s in screens}
    combo_inner = [
        s for s in ctx.shapes
        if id(s) not in screen_ids
        and intersects_region(s, inner_L, inner_T, inner_R, inner_B, 0.1)
        and (s.is_rect_like or s.is_ellipse or s.is_line)
    ]
    combo_ok = len(combo_inner) >= 2

    pictures = [
        s for s in ctx.shapes
        if s.has_blip
        and intersects_region(s, inner_L, inner_T, inner_R, inner_B, 0.1)
    ]
    picture_ok = bool(pictures)

    inner_ok = lines_ok or combo_ok or picture_ok

    hit = bool(screens) and inner_ok
    return result(
        "第1页中部电化学分析仪前屏幕",
        hit,
        (
            f"蓝色显示区 {len(screens)}，"
            f"水平轴 {len(h_axes)}/竖直轴 {len(v_axes)}/折线段 {len(diag_segs)}(+短碎 {len(non_axis_segs)})，"
            f"组合形状 {len(combo_inner)}，图片 {len(pictures)}"
        ),
    )


def check_analyzer_name(ctx: EvaluationContext) -> RuleResult:
    # 细则区域：距左15.4–18.8cm，距上10.5–11.5cm
    REGION_L, REGION_T, REGION_R, REGION_B = 15.4, 10.5, 18.8, 11.5
    texts = find_text_in(ctx, "电化学分析仪", REGION_L, REGION_T, REGION_R, REGION_B, tol=0.15)
    matched = []
    for s in texts:
        # 字体：黑体、微软雅黑或相近中文字体
        font_ok = True
        if s.font_name:
            fn = s.font_name.lower()
            font_ok = any(k in fn for k in (
                "黑体", "微软雅黑", "simhei", "microsoft yahei",
                "等线", "dengxian", "思源黑体", "noto sans sc",
                "heiti", "pingfang",
            ))
        # 字号12–14磅（严格无容差）
        size_ok = s.font_size is not None and in_range(s.font_size, 12.0, 14.0, 0.0)
        # 颜色为白色 —— 明确可解析(schemeClr 已在 first_xml_rgb 中反查为 RGB)
        color_ok = s.text_color is not None and is_whiteish(s.text_color)
        if font_ok and size_ok and color_ok:
            matched.append(s)
    hit = bool(matched)
    return result("第1页中部电化学分析仪名称文本", hit, f"区域内文本 {len(texts)} 个，全条件符合 {len(matched)} 个")



def check_indicator_lights(ctx: EvaluationContext) -> RuleResult:
    # 细则区域：距左18.5–19.5cm，距上10–12cm
    REGION_L, REGION_T, REGION_R, REGION_B = 18.5, 10.0, 19.5, 12.0
    # 直径0.15–0.28cm（宽和高均在此范围内）；椭圆形圆点
    # 注：实际PPT中三个圆点cy分别约9.99/10.53/11.06，略超出细则上边界10cm
    # 用intersects_region并给0.1cm容差以覆盖临界情况
    lights = [
        s for s in ctx.ellipse_shapes
        if intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.1)
        and 0.15 <= s.w_abs <= 0.28
        and 0.15 <= s.h_abs <= 0.28
    ]

    # 三灯竖向排列：按cy排序后判断上/中/下颜色
    lights_sorted = sorted(lights, key=lambda s: s.cy)
    three_lights = lights_sorted[:3] if len(lights_sorted) >= 3 else lights_sorted

    # 上灯绿色、中灯蓝色、下灯红色
    top_green = bool(three_lights) and any(is_green(c) for c in three_lights[0].fill_colors) if len(three_lights) >= 1 else False
    mid_blue = bool(three_lights) and any(is_blue(c) for c in three_lights[1].fill_colors) if len(three_lights) >= 2 else False
    bot_red = bool(three_lights) and any(is_red(c) for c in three_lights[2].fill_colors) if len(three_lights) >= 3 else False

    hit = len(lights) >= 3 and top_green and mid_blue and bot_red
    return result(
        "第1页中部电化学分析仪指示灯组",
        hit,
        f"圆灯 {len(lights)} 个，上绿/中蓝/下红={top_green}/{mid_blue}/{bot_red}",
    )


def check_signal_box(ctx: EvaluationContext) -> RuleResult:
    # 虚线框区域：距左16.3–21.7cm，距上6–7.5cm
    BOX_L, BOX_T, BOX_R, BOX_B = 16.3, 6.0, 21.7, 7.5

    # 方式1：原生蓝色虚线圆角矩形，尺寸4.8–5.2cm × 1.2–1.6cm
    native_boxes = [
        s for s in ctx.rect_shapes
        if center_in(s, BOX_L, BOX_T, BOX_R, BOX_B, 0.3)
        and s.geom == "roundRect"
        and s.dashed
        and is_blue(s.line_color)
        and line_width_in_range(s.line_width, 1.0, 1.5, 0.0)
        and in_range(s.w_abs, 4.8, 5.2, 0.0)
        and in_range(s.h_abs, 1.2, 1.6, 0.0)
        and (not s.fill_colors or any(is_whiteish(c) for c in s.fill_colors))
    ]

    # 方式2：多段蓝色短线组成的虚线边框（用segmented_blue_border已检测，复用即可）
    box_hit, box_segments = segmented_blue_border(ctx, BOX_L, BOX_T, BOX_R, BOX_B, min_segments=14)
    # 进一步验证线段线宽1–1.5磅
    seg_lines = line_in_region(ctx, BOX_L, BOX_T, BOX_R, BOX_B,
                               color_pred=is_blue, width_range=(1.0, 1.5),
                               width_tol=0.0, region_tol=0.1)
    box_width_ok = len(seg_lines) >= 14

    # 文本"信号采集与处理"：在框内，字体、字号14–16磅、颜色深蓝色
    raw_texts = find_text_in(ctx, "信号采集与处理", BOX_L, BOX_T, BOX_R, BOX_B, tol=0.15)
    valid_texts = []
    for s in raw_texts:
        font_ok = True
        if s.font_name:
            fn = s.font_name.lower()
            font_ok = any(k in fn for k in (
                "黑体", "微软雅黑", "simhei", "microsoft yahei",
                "等线", "dengxian", "思源黑体", "noto sans sc",
                "heiti", "pingfang",
            ))
        size_ok = s.font_size is not None and in_range(s.font_size, 14.0, 16.0, 0.0)
        color_ok = s.text_color is not None and is_dark_blue(s.text_color)
        if font_ok and size_ok and color_ok:
            valid_texts.append(s)

    # 箭头区域：距左18–19cm，距上7.3–9cm；蓝色竖向虚线单箭头，线宽1–1.5磅，朝下
    ARROW_L, ARROW_T, ARROW_R, ARROW_B = 18.0, 7.3, 19.0, 9.0

    # 原生虚线箭头
    native_arrows = line_in_region(
        ctx, ARROW_L, ARROW_T, ARROW_R, ARROW_B,
        color_pred=is_blue, width_range=(1.0, 1.5),
        dashed=True, arrow=True, vertical=True,
        width_tol=0.0, region_tol=0.05,
    )
    # 分段虚线：多段蓝色短竖线 + 三角形箭尖
    #   - 短竖线线宽 1.0–1.5 磅(与 rubric 一致, 不再放宽到 1.8);
    #   - 三角尖必须位于短竖线组的下方, 且形状本身"朝下"(downArrow 或 三角形 rotation ≈ 0°/180°);
    blue_vsegs = [
        s for s in ctx.line_shapes
        if intersects_region(s, ARROW_L, ARROW_T, ARROW_R, ARROW_B, 0.05)
        and is_blue(s.line_color)
        and s.w_abs <= 0.05 and s.h_abs >= 0.05
        and line_width_in_range(s.line_width, 1.0, 1.5, 0.0)
    ]
    # 计算竖线组底端 y (若无竖线, 用 ARROW_T 兜底, 便于后续与三角尖比较)
    vseg_bottom_y = max((s.y2 for s in blue_vsegs), default=ARROW_T)

    def tip_points_down(t: ShapeInfo) -> bool:
        """三角尖朝下判定:
          - downArrow 预设几何默认朝下, 视 rot 是否偏转 >45° 决定;
          - triangle 默认尖角朝上, 需 rot≈180°(即 flipV 或 |rot|>=135°) 才朝下;
          - 未取到 rotation 属性时保守放行(避免误伤)。
        """
        r = getattr(t, "rotation", 0.0) or 0.0
        flip_v = getattr(t, "flip_v", False)
        eff = ((r + 180.0) % 360.0) - 180.0  # (-180,180]
        if t.geom == "downArrow":
            # 默认朝下, 允许 ±45° 内的旋转
            return abs(eff) <= 45.0
        if t.geom == "triangle":
            # 默认朝上, 需通过 flipV 或旋转 ≈180° 才朝下
            if flip_v:
                return True
            return abs(eff) >= 135.0
        return False

    raw_tips = [
        s for s in ctx.shapes
        if s.geom in ("triangle", "downArrow")
        and intersects_region(s, ARROW_L, 7.8, ARROW_R, ARROW_B + 0.4, 0.1)
        and (is_blue(s.line_color) or any(is_blue(c) for c in s.fill_colors))
    ]
    arrow_tips = [
        s for s in raw_tips
        if tip_points_down(s)
        # 箭尖位于短竖线组下方(容差 0.35cm): 用尖端 y (三角形/箭头的 y2) 与竖线组底端比较
        and s.y2 >= vseg_bottom_y - 0.35
    ]
    seg_arrow_ok = len(blue_vsegs) >= 2 and bool(arrow_tips)

    # 末端指向仪器顶面（y ≈ 8.5–9.5cm）
    analyzer_top = [
        s for s in ctx.shapes
        if intersects_region(s, 13.0, 8.5, 22.0, 10.0, 0.2)
        and (s.is_rect_like or s.geom == "parallelogram")
        and s.w_abs >= 3.0
    ]
    tip_near_top = False
    for a in native_arrows:
        tip_y = a.y2 if a.box.height >= 0 else a.y1
        if any(abs(tip_y - g.y1) <= 0.6 for g in analyzer_top):
            tip_near_top = True
    for tip in arrow_tips:
        if any(abs(tip.y2 - g.y1) <= 0.6 for g in analyzer_top):
            tip_near_top = True

    box_ok = bool(native_boxes) or (box_hit and box_width_ok)
    arrows_ok = (bool(native_arrows) or seg_arrow_ok) and tip_near_top
    hit = box_ok and bool(valid_texts) and arrows_ok
    return result(
        "第1页中部上方虚线框及向下箭头",
        hit,
        f"原生虚线框 {len(native_boxes)}，分段框线 {box_segments}段，文本 {len(valid_texts)}/{len(raw_texts)}，"
        f"原生箭头 {len(native_arrows)}，分段箭头线{len(blue_vsegs)}/尖{len(arrow_tips)}，末端指向仪器={'是' if tip_near_top else '否'}",
    )



def check_control_box(ctx: EvaluationContext) -> RuleResult:
    box_hit, box_segments = segmented_blue_border(ctx, 14.4, 13.1, 19.4, 15.0, min_segments=12)
    boxes = [s for s in rect_in_region(ctx, 14.4, 13.1, 19.4, 15.0, dashed=True, color_pred=is_blue, min_w=4.2, min_h=1.3) if s.geom == "roundRect" and 4.2 <= s.w_abs <= 4.7 and 1.3 <= s.h_abs <= 1.7]
    texts = [s for s in find_text_in(ctx, ["电位/电流控制", "电位 / 电流控制"], 14.4, 13.1, 19.4, 15.0, tol=0.1) if has_text_style(s, 14, 16, is_dark_blue)]
    arrows = line_in_region(ctx, 16.0, 11.8, 17.6, 13.7, color_pred=is_blue, width_range=(1.0, 1.5), vertical=True, arrow=True, width_tol=0.0, direction="up")
    hit = (bool(boxes) or box_hit) and bool(texts) and bool(arrows)
    return result("第1页中部下方虚线框及向上箭头", hit, f"虚线框短线 {box_segments} 段，文本 {len(texts)}，箭头 {len(arrows)}")


def check_analyzer_to_computer_line(ctx: EvaluationContext) -> RuleResult:
    # 细则区域：距左21.2–24.8cm，距上9.8–11cm
    REGION_L, REGION_T, REGION_R, REGION_B = 21.2, 9.8, 24.8, 11.0

    def is_black_solid_line(s: ShapeInfo) -> bool:
        """黑色单实线，线宽 1.5–2.5 磅。"""
        return (
            (is_blackish(s.line_color) or is_dark(s.line_color))
            and not s.dashed
            and line_width_in_range(s.line_width, 1.5, 2.5, 0.0)
        )

    # 主路径：单条原生线 / 弧形连接符横跨全区域（rubric "黑色弧形或近水平单实线"）
    single_line_candidates: list[ShapeInfo] = [
        s for s in ctx.line_shapes
        if intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.15)
        and is_black_solid_line(s)
        and s.w_abs >= 2.5                # 跨度至少 2.5cm（覆盖分析仪右缘→笔记本左缘的大部分）
        and s.x1 <= 22.5                  # 起点位于区域左段（分析仪右侧）
        and s.x2 >= 24.0                  # 终点位于区域右段（笔记本左侧接口）
        and s.h_abs <= max(1.2, s.w_abs)  # 弧形允许一定高度但整体不高于宽度，且不超出区域高度 1.2cm
    ]
    single_line_ok = bool(single_line_candidates)

    # 兼容分支：多段近水平短线拼接。段数门槛下调至 2，避免误判合规单线
    line_segs: list[ShapeInfo] = [
        s for s in ctx.line_shapes
        if intersects_region(s, REGION_L, REGION_T, REGION_R, REGION_B, 0.1)
        and is_black_solid_line(s)
        and s.h_abs <= max(0.08, s.w_abs * 0.35)
    ]
    start_near_analyzer = [
        s for s in line_segs
        if intersects_region(s, 21.2, REGION_T, 22.0, REGION_B, 0.15)
    ]
    end_near_computer = [
        s for s in line_segs
        if intersects_region(s, 24.0, REGION_T, 24.8, REGION_B, 0.15)
    ]
    multi_span_ok = bool(line_segs) and (max(s.x2 for s in line_segs) - min(s.x1 for s in line_segs) >= 3.0)
    multi_ok = (
        len(line_segs) >= 2
        and bool(start_near_analyzer)
        and bool(end_near_computer)
        and multi_span_ok
    )

    # 笔记本电脑左侧接口块：仅作为辅助信号，不作硬约束（单实线主路径不依赖接口块）
    computer_ports = [
        s for s in ctx.shapes
        if intersects_region(s, 24.2, 9.7, 25.1, 10.6, 0.1)
        and (s.is_rect_like or s.is_line)
        and (is_blackish(s.line_color) or is_dark(s.line_color) or any(is_blackish(c) or is_dark(c) for c in s.fill_colors))
    ]

    hit = single_line_ok or multi_ok
    return result(
        "第1页中部右侧连接线",
        hit,
        f"单实线/弧形 {len(single_line_candidates)}，多段线 {len(line_segs)}（起点 {len(start_near_analyzer)}，终点 {len(end_near_computer)}，跨度合规={'是' if multi_span_ok else '否'}），接口块 {len(computer_ports)}",
    )


def check_laptop_outline(ctx: EvaluationContext) -> RuleResult:
    # 细则区域：距左22.5–32cm，距上6.2–12.5cm
    REGION_L, REGION_T, REGION_R, REGION_B = 22.5, 6.2, 32.0, 12.5

    def is_deep_blue_or_bluegrey_or_blueblack(c: tuple[int, int, int] | None) -> bool:
        """机身颜色：深蓝色、深灰蓝色或蓝黑色 —— 只要属于蓝色/黑色/深色/灰色家族即视为合规"""
        return is_blue(c) or is_blackish(c) or is_dark(c) or is_grayish(c)

    def has_laptop_outline(s: ShapeInfo) -> bool:
        """边线深蓝或黑色单实线，线宽0.75–1.25磅（未显式设置按默认1.0磅处理）"""
        is_dark_line = (
            is_blue(s.line_color)
            or is_blackish(s.line_color)
            or is_dark(s.line_color)
            or is_grayish(s.line_color)
        )
        return (
            not s.dashed
            and is_dark_line
            and line_width_in_range(s.line_width, 0.75, 1.25, 0.0)
        )

    # 屏幕（包含屏幕框体）：在区域上部，大面积矩形/圆角矩形，机身深蓝/蓝黑色填充，边线合规
    screen_parts = [
        s for s in ctx.rect_shapes
        if intersects_region(s, REGION_L, REGION_T, REGION_R, 12.0, 0.2)
        and s.w_abs >= 5.0 and s.h_abs >= 3.0
        and any(is_deep_blue_or_bluegrey_or_blueblack(c) for c in s.fill_colors)
        and has_laptop_outline(s)
    ]

    # 键盘底座：在区域下部（y>11cm），较宽扁矩形或梯形，机身深色填充，边线合规
    base_parts = [
        s for s in ctx.shapes
        if intersects_region(s, REGION_L, 11.0, REGION_R, REGION_B, 0.2)
        and (s.is_rect_like or s.geom == "trapezoid")
        and s.w_abs >= 5.0 and s.h_abs <= 2.5
        and any(is_deep_blue_or_bluegrey_or_blueblack(c) for c in s.fill_colors)
    ]

    # 后侧折叠连接部（合页处）：介于屏幕底部和底座之间，小宽矩形/弧形，深色
    hinge_parts = [
        s for s in ctx.rect_shapes
        if intersects_region(s, REGION_L, 11.5, REGION_R, 13.0, 0.2)
        and s.w_abs >= 3.0 and s.h_abs <= 1.5
        and any(is_deep_blue_or_bluegrey_or_blueblack(c) for c in s.fill_colors)
    ]

    hit = bool(screen_parts) and bool(base_parts) and bool(hinge_parts)
    return result(
        "第1页右侧笔记本电脑外轮廓",
        hit,
        f"屏幕框体 {len(screen_parts)}，键盘底座 {len(base_parts)}，折叠连接 {len(hinge_parts)}",
    )


def check_laptop_screen_charts(ctx: EvaluationContext) -> RuleResult:
    # 屏幕整体区域：距左22.5–31.8cm，距上6.4–11cm
    SCR_L, SCR_T, SCR_R, SCR_B = 22.5, 6.4, 31.8, 11.0

    # 屏幕：白色或极浅灰白填充，深蓝色边框，尺寸足够大
    screen = [
        s for s in ctx.rect_shapes
        if intersects_region(s, SCR_L, SCR_T, SCR_R, SCR_B, 0.15)
        and s.w_abs >= 6.0 and s.h_abs >= 3.0
        and any(is_whiteish(c) or is_light_grayish(c) for c in s.fill_colors)
        and is_dark_blue(s.line_color)
        and not s.dashed
    ]

    # 屏幕内上方浅灰色界面横条：浅灰/白色填充，宽>=5cm，高<=0.7cm
    ui_bar = [
        s for s in ctx.rect_shapes
        if intersects_region(s, SCR_L, SCR_T, SCR_R, 7.8, 0.1)
        and s.w_abs >= 5.0 and s.h_abs <= 0.7
        and any(is_light_grayish(c) or is_whiteish(c) for c in s.fill_colors)
    ]

    # 小圆点装饰：在上方横条左侧区域，小椭圆，直径0.08–0.25cm
    ui_dots = [
        s for s in ctx.ellipse_shapes
        if intersects_region(s, 22.7, SCR_T, 24.5, 7.6, 0.1)
        and 0.08 <= s.w_abs <= 0.25 and 0.08 <= s.h_abs <= 0.25
        and any(is_light_grayish(c) or is_whiteish(c) or is_blue(c) for c in s.fill_colors)
    ]

    # 左半屏幕蓝色折线图：距左25–28 × 7.3–10.5cm，蓝色曲线线宽0.75–1.25磅
    CHART_L_L, CHART_T, CHART_L_R, CHART_B = 25.0, 7.3, 28.0, 10.5
    blue_lines = line_in_region(
        ctx, CHART_L_L, CHART_T, CHART_L_R, CHART_B,
        color_pred=is_blue, width_range=(0.75, 1.25),
        min_len=0.15, width_tol=0.0, region_tol=0.05,
    )
    # 左半浅灰色坐标轴，线宽0.75–1.25磅
    left_axes = line_in_region(
        ctx, CHART_L_L, CHART_T, CHART_L_R, CHART_B,
        color_pred=is_light_grayish, width_range=(0.75, 1.25),
        min_len=0.3, width_tol=0.0, region_tol=0.05,
    )

    # 右半屏幕红色分布曲线图：距左27.8–30.2 × 7.3–10.5cm，红色曲线线宽0.75–1.25磅
    CHART_R_L, CHART_R_R = 27.8, 30.2
    red_lines = line_in_region(
        ctx, CHART_R_L, CHART_T, CHART_R_R, CHART_B,
        color_pred=is_red, width_range=(0.75, 1.25),
        min_len=0.15, width_tol=0.0, region_tol=0.05,
    )
    # 右半浅灰色坐标轴，线宽0.75–1.25磅
    right_axes = line_in_region(
        ctx, CHART_R_L, CHART_T, CHART_R_R, CHART_B,
        color_pred=is_light_grayish, width_range=(0.75, 1.25),
        min_len=0.3, width_tol=0.0, region_tol=0.05,
    )

    hit = (
        bool(screen)
        and bool(ui_bar)
        and len(ui_dots) >= 2
        and len(blue_lines) >= 2
        and len(left_axes) >= 2
        and len(red_lines) >= 2
        and len(right_axes) >= 2
    )
    return result(
        "第1页右侧笔记本电脑屏幕区域",
        hit,
        f"屏幕 {len(screen)}，界面横条 {len(ui_bar)}，圆点 {len(ui_dots)}，"
        f"蓝线 {len(blue_lines)}/左轴 {len(left_axes)}，红线 {len(red_lines)}/右轴 {len(right_axes)}",
    )



def check_computer_text(ctx: EvaluationContext) -> RuleResult:
    # 细则：位于笔记本电脑下方中央，距左25.0–27.4cm、距上12.4–13.8cm
    # 用 tol=0.3 允许文本框中心在给定范围内有合理偏移
    texts = find_text_in(ctx, "计算机", 25.0, 12.4, 27.4, 13.8, tol=0.3)

    def is_black_or_dark_grey(c: tuple[int, int, int] | None) -> bool:
        """黑色或深灰色 —— 只要属于黑色/深色/灰色家族即视为合规"""
        return is_blackish(c) or is_dark(c) or is_grayish(c)

    # 细则：字体黑体/微软雅黑或相近中文字体；字号13–15磅；颜色黑色或深灰色
    # size_tol=0.0 严格匹配13–15磅；字体通过 has_text_style 内 SANS_FONTS 覆盖
    hit = any(has_text_style(s, 13, 15, is_black_or_dark_grey, size_tol=0.0) for s in texts)
    return result("第1页右侧笔记本电脑下方文本", hit, f"指定区域文本 {len(texts)} 个")


def check_data_box(ctx: EvaluationContext) -> RuleResult:
    # 细则：蓝色虚线圆角矩形，位于距左25.2–31.4cm、距上3.4–5.2cm
    # 线宽1–1.5磅，填充白色或无填充，宽5.2–5.6cm，高1.3–1.7cm
    _, box_segments = segmented_blue_border(ctx, 25.2, 3.4, 31.4, 5.2, min_segments=14)
    boxes = [
        s for s in rect_in_region(ctx, 25.2, 3.4, 31.4, 5.2, dashed=True, color_pred=is_blue, min_w=5.2, min_h=1.3)
        if s.geom == "roundRect"
        and 5.2 <= s.w_abs <= 5.6
        and 1.3 <= s.h_abs <= 1.7
        and line_width_in_range(s.line_width, 1.0, 1.5, 0.0)  # 细则：线宽1–1.5磅（未显式设置按默认1.0磅处理）
        and (not s.fill_colors or any(is_whiteish(c) for c in s.fill_colors))  # 细则：填充白色或无填充
    ]
    box_hit = bool(boxes)

    # 细则：内部文本"数据分析与可视化"，字体黑体/微软雅黑等，字号18–24磅，颜色深蓝色
    texts = [
        s for s in find_text_in(ctx, "数据分析与可视化", 25.2, 3.4, 31.4, 5.2, tol=0.3)
        if has_text_style(s, 18, 24, is_dark_blue, size_tol=0.0)
    ]

    # 细则：蓝色竖向虚线单箭头，线宽1–1.5磅，箭头朝下
    # 位于距左27.5–28.8cm、距上4.8–6.7cm
    arrows = line_in_region(
        ctx, 27.5, 4.8, 28.8, 6.7,
        color_pred=is_blue,
        width_range=(1.0, 1.5),
        dashed=True,
        arrow=True,
        vertical=True,
        width_tol=0.0,
        region_tol=0.3,
        direction="down",
    )

    hit = box_hit and bool(texts) and bool(arrows)
    return result("第1页右上虚线框及向下箭头", hit, f"虚线框短线 {box_segments} 段，原生框 {len(boxes)}，文本 {len(texts)}，向下虚线箭头 {len(arrows)}")


def check_legend_frame(ctx: EvaluationContext) -> RuleResult:
    # 细则：蓝色虚线圆角矩形，位于距左3.6–30.3cm、距上15.3–18.2cm
    # 宽26–27cm，高2.2–2.5cm，填充白色或无填充，边线蓝色虚线，线宽1–1.5磅
    _, box_segments = segmented_blue_border(ctx, 3.6, 15.3, 30.3, 18.2, min_segments=40)
    boxes = [
        s for s in rect_in_region(ctx, 3.6, 15.3, 30.3, 18.2, dashed=True, color_pred=is_blue, min_w=26.0, min_h=2.2)
        if s.geom == "roundRect"
        and within_region(s, 3.6, 15.3, 30.3, 18.2, 0.3)  # 细则：位于指定范围内
        and 26.0 <= s.w_abs <= 27.0          # 细则：宽26–27cm
        and 2.2 <= s.h_abs <= 2.5            # 细则：高2.2–2.5cm
        and line_width_in_range(s.line_width, 1.0, 1.5, 0.0)  # 细则：线宽1–1.5磅（未显式设置按默认1.0磅处理）
        and (not s.fill_colors or any(is_whiteish(c) for c in s.fill_colors))  # 细则：填充白色或无填充
    ]
    hit = bool(boxes)
    return result("第1页底部图例外框", hit, f"底部蓝色虚线短线 {box_segments} 段，原生虚线框 {len(boxes)} 个")


def check_legend_black(ctx: EvaluationContext) -> RuleResult:
    # 细则：黑色水平单实线，位于距左4.4–6.1cm、距上16.0–16.7cm
    # 长度0.9–1.5cm，线宽1.5–2.5磅
    lines = [
        s for s in line_in_region(
            ctx, 4.4, 16.0, 6.1, 16.7,
            color_pred=is_blackish,
            width_range=(1.5, 2.5),
            dashed=False,
            arrow=False,
            horizontal=True,
            min_len=0.9,
            max_len=1.5,
            width_tol=0.0,
            region_tol=0.2,
        )
        if within_region(s, 4.4, 16.0, 6.1, 16.7, 0.2)
    ]

    def is_black_or_dark_grey(c: tuple[int, int, int] | None) -> bool:
        """黑色或深灰色 —— 只要属于黑色/深色/灰色家族即视为合规。"""
        return is_blackish(c) or is_dark(c) or is_grayish(c)

    # 细则：文本位于黑色线段右侧，距左6.0–12.2cm、距上16.0–16.7cm
    # 内容为“工作电极连接线（信号输出）”，字体黑体/微软雅黑等，字号11–13磅，颜色黑色或深灰色
    texts = [
        s for s in find_text_in(ctx, "工作电极连接线（信号输出）", 6.0, 16.0, 12.2, 16.7, tol=0.2)
        if has_text_style(s, 11, 13, is_black_or_dark_grey, size_tol=0.0)
    ]

    hit = bool(lines) and bool(texts)
    return result("第1页底部图例黑色线段", hit, f"黑色水平单实线 {len(lines)}，文本 {len(texts)}")


def check_legend_green(ctx: EvaluationContext) -> RuleResult:
    # 细则：绿色水平单实线，位于距左4.4–6.1cm、距上16.5–17.8cm
    # 长度0.9–1.5cm，线宽1.5–2.5磅
    lines = [
        s for s in line_in_region(
            ctx, 4.4, 16.5, 6.1, 17.8,
            color_pred=is_green,
            width_range=(1.5, 2.5),
            dashed=False,
            arrow=False,
            horizontal=True,
            min_len=0.9,
            max_len=1.5,
            width_tol=0.0,
            region_tol=0.2,
        )
        if within_region(s, 4.4, 16.5, 6.1, 17.8, 0.2)
    ]

    def is_black_or_dark_grey(c: tuple[int, int, int] | None) -> bool:
        """黑色或深灰色 —— 只要属于黑色/深色/灰色家族即视为合规。"""
        return is_blackish(c) or is_dark(c) or is_grayish(c)

    # 细则：文本位于绿色线段右侧，距左6.0–12.2cm、距上16.5–17.8cm
    # 内容为"参比电极连接线（电位参照）"，字体黑体/微软雅黑等，字号11–13磅，颜色黑色或深灰色
    texts = [
        s for s in find_text_in(
            ctx, ["参比电极连接线（电位参照）"],
            6.0, 16.5, 12.2, 17.8, tol=0.2,
        )
        if has_text_style(s, 11, 13, is_black_or_dark_grey, size_tol=0.0)
    ]

    hit = bool(lines) and bool(texts)
    return result("第1页底部图例绿色线段", hit, f"绿色水平单实线 {len(lines)}，文本 {len(texts)}")


def check_legend_red(ctx: EvaluationContext) -> RuleResult:
    # 细则：红色水平单实线，位于距左14.0–16.0cm、距上16.0–17.0cm
    # 长度0.9–1.5cm，线宽1.5–2.5磅
    lines = [
        s for s in line_in_region(
            ctx, 14.0, 16.0, 16.0, 17.0,
            color_pred=is_red,
            width_range=(1.5, 2.5),
            dashed=False,
            arrow=False,
            horizontal=True,
            min_len=0.9,
            max_len=1.5,
            width_tol=0.0,
            region_tol=0.2,
        )
        if within_region(s, 14.0, 16.0, 16.0, 17.0, 0.2)
    ]

    def is_black_or_dark_grey(c: tuple[int, int, int] | None) -> bool:
        """黑色或深灰色 —— 只要属于黑色/深色/灰色家族即视为合规。"""
        return is_blackish(c) or is_dark(c) or is_grayish(c)

    # 细则：文本位于红色线段右侧，距左15.8–21.4cm、距上16.0–17.0cm
    # 内容为"对电极连接线（电流通路）"，字体黑体/微软雅黑等，字号11–13磅，颜色黑色或深灰色
    texts = [
        s for s in find_text_in(ctx, "对电极连接线（电流通路）", 15.8, 16.0, 21.4, 17.0, tol=0.2)
        if has_text_style(s, 11, 13, is_black_or_dark_grey, size_tol=0.0)
    ]

    hit = bool(lines) and bool(texts)
    return result("第1页底部图例红色线段", hit, f"红色水平单实线 {len(lines)}，文本 {len(texts)}")


def check_legend_blue_arrow(ctx: EvaluationContext) -> RuleResult:
    # 细则：蓝色水平虚线单箭头，位于距左23.2–25.0cm、距上16.4–17.1cm
    # 线宽1–1.5磅，箭头朝右
    arrows = [
        s for s in line_in_region(
            ctx, 23.2, 16.4, 25.0, 17.1,
            color_pred=is_blue,
            width_range=(1.0, 1.5),
            dashed=True,
            arrow=True,
            horizontal=True,
            width_tol=0.0,
            region_tol=0.2,
            direction="right",
        )
        if within_region(s, 23.2, 16.4, 25.0, 17.1, 0.2)
    ]

    def is_black_or_dark_grey(c: tuple[int, int, int] | None) -> bool:
        """黑色或深灰色 —— 只要属于黑色/深色/灰色家族即视为合规。"""
        return is_blackish(c) or is_dark(c) or is_grayish(c)

    # 细则：文本位于蓝色虚线箭头右侧，距左24.9–29.3cm、距上16.4–17.1cm
    # 内容为"信号/控制/数据流向"，字体黑体/微软雅黑等，字号11–13磅，颜色黑色或深灰色
    texts = [
        s for s in find_text_in(ctx, "信号/控制/数据流向", 24.9, 16.4, 29.3, 17.1, tol=0.2)
        if has_text_style(s, 11, 13, is_black_or_dark_grey, size_tol=0.0)
    ]

    hit = bool(arrows) and bool(texts)
    return result("第1页底部图例蓝色虚线箭头", hit, f"蓝色水平虚线单箭头 {len(arrows)}，文本 {len(texts)}")


def check_overall_layout(ctx: EvaluationContext) -> RuleResult:
    # 细则：以下各区域均须在页面可视范围内有对象
    # 1. 外侧蓝色边框 — 贴近页面四边
    PAGE_W = ctx.slide_width
    PAGE_H = ctx.slide_height

    # 外侧蓝色边框
    border_in_view = any(
        s for s in ctx.rect_shapes
        if s.x1 <= 0.4 and s.y1 <= 0.4
        and s.x2 >= PAGE_W - 0.4 and s.y2 >= PAGE_H - 0.4
    )

    # 左侧烧杯与三根电极（左区有足够对象）
    left_objects = count_shapes((3.0, 3.5, 12.5, 14.5), ctx.shapes)
    left_ok = left_objects >= 15

    # 三条彩色连接线（黑/绿/红，横跨中左区域）
    black_lines = line_in_region(ctx, 7.0, 3.5, 14.0, 7.5, color_pred=lambda c: is_blackish(c) or is_dark(c), min_len=1.0)
    green_lines = line_in_region(ctx, 7.0, 3.5, 14.0, 12.5, color_pred=is_green, min_len=1.0)
    red_lines   = line_in_region(ctx, 7.0, 3.5, 14.0, 12.5, color_pred=is_red,   min_len=1.0)
    connection_lines_ok = bool(black_lines) and bool(green_lines) and bool(red_lines)

    # 中部电化学分析仪
    center_objects = count_shapes((13.0, 5.5, 22.5, 15.5), ctx.shapes)
    center_ok = center_objects >= 15

    # 右侧笔记本电脑
    right_objects = count_shapes((22.5, 3.0, 32.5, 14.5), ctx.shapes)
    right_ok = right_objects >= 15

    # 顶部与中部三个蓝色虚线标注框（顶右、中上、中下）
    blue_dashed_boxes = [
        s for s in ctx.rect_shapes
        if s.dashed and is_blue(s.line_color)
        and intersects_region(s, 3.0, 3.0, 32.0, 15.5, 0.3)
    ]
    annotation_boxes_ok = len(blue_dashed_boxes) >= 3

    # 底部图例框及四组图例均在页面内
    legend_objects = count_shapes((3.4, 15.0, 30.8, 18.5), ctx.shapes)
    legend_box_in_view = any(
        s for s in ctx.rect_shapes
        if s.dashed and is_blue(s.line_color)
        and intersects_region(s, 3.0, 14.5, 31.5, 18.5, 0.3)
        and s.w_abs >= 20.0
    )
    legend_ok = legend_objects >= 8 and legend_box_in_view

    # 所有对象不超出页面可视范围
    in_view_ok = not any(
        s for s in ctx.shapes
        if s.x2 < -0.3 or s.y2 < -0.3
        or s.x1 > PAGE_W + 0.3 or s.y1 > PAGE_H + 0.3
    )

    # 文本未超出对应框体（文本中心须在页面内）
    text_overflow = [
        s for s in ctx.text_shapes
        if s.cx < 0 or s.cy < 0 or s.cx > PAGE_W or s.cy > PAGE_H
    ]
    text_in_view_ok = len(text_overflow) == 0

    # 主要图形之间没有明显重叠导致无法阅读
    # （文本框与文本框之间重叠面积超过较小框50%视为严重重叠）
    text_overlaps = 0
    texts = ctx.text_shapes
    for i, a in enumerate(texts):
        for b in texts[i + 1:]:
            if overlap_area(a.box, b.box) > min(a.box.area, b.box.area) * 0.50:
                text_overlaps += 1
    no_bad_overlap = text_overlaps <= 2

    hit = (
        border_in_view
        and left_ok
        and connection_lines_ok
        and center_ok
        and right_ok
        and annotation_boxes_ok
        and legend_ok
        and in_view_ok
        and text_in_view_ok
        and no_bad_overlap
    )
    return result(
        "第1页整体排版",
        hit,
        (
            f"蓝框={'有' if border_in_view else '无'}，"
            f"左区{left_objects}个，连接线(黑{len(black_lines)}/绿{len(green_lines)}/红{len(red_lines)})，"
            f"中区{center_objects}个，右区{right_objects}个，"
            f"蓝色虚线标注框{len(blue_dashed_boxes)}个，"
            f"图例区{legend_objects}个(框={'有' if legend_box_in_view else '无'})，"
            f"越界={'无' if in_view_ok else '有'}，文本越框={'无' if text_in_view_ok else f'{len(text_overflow)}处'}，"
            f"文本严重重叠{text_overlaps}处"
        ),
    )


def dimension2_items() -> list[ScoreItem]:
    return [
        ScoreItem(1, "第1页白色页面背景", check_background),
        ScoreItem(1, "第1页外侧蓝色矩形边框", check_outer_border),
        ScoreItem(1, "第1页左侧烧杯外轮廓", check_beaker_outline),
        ScoreItem(1, "第1页左侧烧杯顶部开口椭圆", check_beaker_top_ellipse),
        ScoreItem(3, "第1页左侧烧杯内部液面与液体填充", check_liquid_and_bubbles),
        ScoreItem(1, "第1页左侧烧杯底部文本", check_sample_text),
        ScoreItem(3, "第1页左侧绿色参比电极", check_reference_electrode),
        ScoreItem(3, "第1页左侧中间白蓝工作电极", check_working_electrode),
        ScoreItem(3, "第1页右侧红色对电极", check_counter_electrode),
        ScoreItem(1, '第1页左侧"参比电极"文本', check_reference_label_arrow),
        ScoreItem(1, '第1页中部"工作电极"文本', check_working_label),
        ScoreItem(1, '第1页右侧"对电极"文本', check_counter_label_arrow),
        ScoreItem(1, '第1页顶部"分析物输入"文本', check_analyte_input),
        ScoreItem(3, "第1页三条电极连接线", check_electrode_connection_lines),
        ScoreItem(5, "第1页中部电化学分析仪外壳", check_analyzer_shell),
        ScoreItem(3, "第1页中部电化学分析仪前屏幕", check_analyzer_screen),
        ScoreItem(1, "第1页中部电化学分析仪名称文本", check_analyzer_name),
        ScoreItem(1, "第1页中部电化学分析仪指示灯组", check_indicator_lights),
        ScoreItem(1, "第1页中部上方虚线框及向下箭头", check_signal_box),
        ScoreItem(1, "第1页中部下方虚线框及向上箭头", check_control_box),
        ScoreItem(1, "第1页中部右侧连接线", check_analyzer_to_computer_line),
        ScoreItem(5, "第1页右侧笔记本电脑外轮廓", check_laptop_outline),
        ScoreItem(5, "第1页右侧笔记本电脑屏幕区域", check_laptop_screen_charts),
        ScoreItem(1, "第1页右侧笔记本电脑下方文本", check_computer_text),
        ScoreItem(1, "第1页右上虚线框及向下箭头", check_data_box),
        ScoreItem(1, "第1页底部图例外框", check_legend_frame),
        ScoreItem(1, "第1页底部图例黑色线段", check_legend_black),
        ScoreItem(1, "第1页底部图例绿色线段", check_legend_green),
        ScoreItem(1, "第1页底部图例红色线段", check_legend_red),
        ScoreItem(1, "第1页底部图例蓝色虚线箭头", check_legend_blue_arrow),
        ScoreItem(3, "第1页整体排版", check_overall_layout),
    ]


def run_dimension2(ctx: EvaluationContext) -> list[RuleResult]:
    results = []
    for item in dimension2_items():
        try:
            rr = item.checker(ctx)
            rr.label = item.label
            rr.points = item.points
        except Exception as exc:
            rr = RuleResult(
                item.label,
                False,
                f"检测异常：{type(exc).__name__}: {exc}",
                item.points,
                errored=True,
            )
        results.append(rr)
    return results


SCRIPT_ID = "053"

# 被评估文档的扩展名 —— 仅支持 .pptx (python-pptx + xml.etree + zipfile 静态解析)
SUPPORTED_EXTS = (".pptx",)


def _locate_document(dir_path: Path) -> Optional[Path]:
    """在脚本所在目录内定位待评估的 .pptx 文件。"""
    # 优先使用与默认文件名一致的目标文件
    preferred = dir_path / DEFAULT_FILE
    if preferred.is_file() and preferred.suffix.lower() == ".pptx":
        return preferred
    # 否则扫描目录，取第一个 .pptx 文件
    candidates: list[Path] = []
    for ext in SUPPORTED_EXTS:
        candidates.extend(sorted(p for p in dir_path.glob(f"*{ext}") if p.is_file()))
    return candidates[0] if candidates else None


def _build_dim2_items(d2: Optional[list[RuleResult]], all_items: list[ScoreItem]) -> list[dict]:
    """将维度2逐项检测结果转为统一约定的字典列表；未跑维度2时全部记为未命中。

    ``detail`` 字段一律返回空字符串（不暴露内部检测细节或异常信息），
    评分与命中判定仍按原有逻辑基于 ``RuleResult`` 计算，不受影响。
    """
    if d2 is None:
        return [
            {
                "rule": item.label,
                "max_delta": item.points,
                "delta": 0,
                "hit": False,
                "detail": "",
            }
            for item in all_items
        ]
    items: list[dict] = []
    for r in d2:
        items.append(
            {
                "rule": r.label,
                "max_delta": r.points,
                "delta": r.points if r.hit else 0,
                "hit": bool(r.hit),
                "detail": "",
            }
        )
    return items


def evaluate(dir_path: str) -> dict:
    """统一入口：接收脚本所在目录路径，脚本自行在该目录内定位并评估 PPT。"""
    all_items = dimension2_items()
    max_score = sum(item.points for item in all_items)
    result_base = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": max_score,
    }

    try:
        base_dir = Path(dir_path)
        if not base_dir.is_dir():
            result_base["status"] = "error"
            result_base["error"] = f"目录不存在或不是文件夹：{dir_path}"
            result_base["dim2_items"] = _build_dim2_items(None, all_items)
            return result_base

        doc_path = _locate_document(base_dir)
        if doc_path is None:
            result_base["status"] = "error"
            result_base["error"] = f"目录中未找到 .pptx 文件：{dir_path}"
            result_base["dim2_items"] = _build_dim2_items(None, all_items)
            return result_base

        result_base["file_name"] = doc_path.name

        ctx = build_context(doc_path)
        d1 = dimension1_checks(ctx)
        dim1_pass = all(r.hit for r in d1)
        result_base["dim1_pass"] = dim1_pass
        if not dim1_pass:
            failed = [r for r in d1 if not r.hit]
            result_base["dim1_reason"] = "；".join(f"{r.label}：{r.detail}" for r in failed)
            result_base["dim2_items"] = _build_dim2_items(None, all_items)
            result_base["total_score"] = 0
            return result_base

        d2 = run_dimension2(ctx)
        result_base["dim2_items"] = _build_dim2_items(d2, all_items)
        result_base["total_score"] = sum(r.points for r in d2 if r.hit)
        return result_base
    except Exception as exc:
        result_base["status"] = "error"
        result_base["error"] = f"{type(exc).__name__}: {exc}"
        if not result_base["dim2_items"]:
            result_base["dim2_items"] = _build_dim2_items(None, all_items)
        return result_base


if __name__ == "__main__":
    # 本地自测：默认使用脚本所在目录；也可通过命令行参数传入其它目录
    debug_dir = sys.argv[1] if len(sys.argv) > 1 else str(Path(__file__).parent)
    print(json.dumps(evaluate(debug_dir), ensure_ascii=False, indent=2))
