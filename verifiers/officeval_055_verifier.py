#!/usr/bin/env python3
"""Automatically evaluate website_outage_editable_flow.pptx against the rubric."""

from __future__ import annotations

import json

SCRIPT_ID = "055"
import math
import os
import re
import sys
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Iterable, Optional
from xml.etree import ElementTree as ET

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:  # pragma: no cover - handled at runtime
    Presentation = None
    MSO_SHAPE_TYPE = None

EMU_PER_CM = 360000
EMU_PER_INCH = 914400
EMU_PER_PT = 12700
CM_TOL = 0.18
SPEC_CM_TOL = 0.03
PT_TOL = 1.1
# PowerPoint's default outline width when a shape/line does not set <a:ln w="…">.
# The rubric measures the *rendered* stroke, so an unset width should be judged
# as if it were the default (1.0pt) rather than silently passing or failing.
DEFAULT_LINE_WIDTH_PT = 1.0

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

FINISH_BODY = (
    "The site was back up and stable within 10 minutes, order submissions resumed, "
    "and a post-incident review has been scheduled for 4:00 PM this afternoon."
)

STAGE_SPECS = [
    {
        "num": "1",
        "title": "start",
        "body": "At 8:40 a.m., the team noticed that the online store was no longer loading for customers.",
        "top_range": (3.4, 4.0),
        "icon_hint": ["warning", "computer", "monitor"],
        # +5 细则逐点核验（strict）；细则未提及的不额外约束。
        "strict": True,
        "body_fill": "white",       # 说明框白色填充
        # 尺寸区间（rubric 修订版）
        "num_size":   (0.8, 1.5),   # 编号圆宽高 0.8-1.5cm
        "title_size": (3.5, 4.6, 0.7, 1.3),  # 标题条 宽 3.5-4.6cm 高 0.7-1.3cm
        "body_size":  (4.5, 5.7, 4.0, 5.0),  # 说明框 宽 4.5-5.7cm 高 4-5cm
        "icon_size":  (2.9, 3.6),   # 浅蓝色圆形宽高 2.9-3.6cm
    },
    {
        "num": "2",
        "title": "check",
        "body": "The support lead checked alerts, confirmed the outage, and opened an incident channel for developers and operations.",
        "top_range": None,
        "icon_hint": ["magnifier", "search", "exclamation", "alert"],
        # +5 细则逐点核验（strict）；细则未提及的不额外约束。
        "strict": True,
        "body_fill": "light_blue",  # 说明框浅蓝色填充
        "right_of_title": "start",  # 位于第 1 阶段右侧
        "icon_orange_badge": True,  # 圆形内含橙色感叹号小圆标
        # 尺寸区间（rubric 修订版）
        "num_size":   (0.5, 1.5),   # 编号圆宽高 0.5-1.5cm
        "title_size": (3.5, 4.7, 0.5, 1.5),  # 标题条 宽 3.5-4.7cm 高 0.5-1.5cm
        "body_size":  (4.5, 5.7, 3.8, 5.0),  # 说明框 宽 4.5-5.7cm 高 3.8-5cm
        "icon_size":  (2.9, 4.0),   # 浅蓝色圆形宽高 2.9-4cm
    },
    {
        "num": "3",
        "title": "problem peak",
        "body": "Orders failed repeatedly because the database connection pool was exhausted during the launch traffic surge.",
        "top_range": None,
        "icon_hint": ["server", "warning", "alert"],
        # +5 细则逐点核验（strict）；细则未提及的不额外约束。
        "strict": True,
        "body_fill": "light_blue",  # 说明框浅蓝色填充
        "center_of_page": True,     # 位于页面中部
        # 尺寸区间（rubric 修订版）
        "num_size":   (0.8, 1.5),   # 编号圆宽高 0.8-1.5cm
        "title_size": (3.5, 4.7, 0.5, 1.5),  # 标题条 宽 3.5-4.7cm 高 0.5-1.5cm
        "body_size":  (4.5, 5.7, 3.8, 4.8),  # 说明框 宽 4.5-5.7cm 高 3.8-4.8cm
        "icon_size":  (2.8, 3.8),   # 浅蓝色圆形宽高 2.8-3.8cm
    },
    {
        "num": "4",
        "title": "solution point",
        "body": "The team rolled back the latest configuration and switched traffic to a healthy backup instance.",
        "top_range": None,
        "icon_hint": ["cloud", "check"],
        # +5 细则逐点核验（strict）；细则未提及的不额外约束。
        "strict": True,
        "body_fill": "light_blue",   # 说明框浅蓝色填充
        "right_of_title": "problem peak",  # 位于第 3 阶段右侧
        # 尺寸区间（rubric 修订版）
        "num_size":   (0.8, 1.5),   # 编号圆宽高 0.8-1.5cm
        "title_size": (3.5, 4.5, 0.5, 1.5),  # 标题条 宽 3.5-4.5cm 高 0.5-1.5cm
        "body_size":  (4.5, 5.7, 4.0, 5.0),  # 说明框 宽 4.5-5.7cm 高 4-5cm
        "icon_size":  (2.8, 3.8),   # 浅蓝色圆形宽高 2.8-3.8cm
    },
    {
        "num": "5",
        "title": "finish",
        "body": FINISH_BODY,
        "top_range": None,
        "icon_hint": ["monitor", "display", "check"],
        # +5 细则逐点核验（strict）；细则未提及的不额外约束。
        "strict": True,
        "body_fill": "light_blue",   # 说明框浅蓝色填充
        "right_side_of_page": True,  # 位于页面右侧
        "num_white": False,          # 细则仅要求深绿色编号圆，未提"白色编号"
        # 尺寸区间（rubric 修订版）
        "num_size":   (0.8, 1.5),   # 编号圆宽高 0.8-1.5cm
        "title_size": (3.5, 4.7, 0.5, 1.5),  # 标题条 宽 3.5-4.7cm 高 0.5-1.5cm
        "body_size":  (5.5, 5.7, 4.0, 5.0),  # 说明框 宽 5.5-5.7cm 高 4-5cm
        "icon_size":  (2.8, 3.8),   # 浅蓝色圆形宽高 2.8-3.8cm
    },
]

EMOTION_WORDS = ["concerned", "tense and focused", "hopeful", "relieved"]
THEME_TEXT = "theme: managing a website outage before a product launch"


@dataclass
class RunInfo:
    text: str = ""
    font_name: Optional[str] = None
    size_pt: Optional[float] = None
    color: Optional[str] = None
    bold: Optional[bool] = None
    italic: Optional[bool] = None


@dataclass
class ShapeInfo:
    name: str
    shape_type: str
    text: str
    left_cm: float
    top_cm: float
    width_cm: float
    height_cm: float
    fill_color: Optional[str] = None
    line_color: Optional[str] = None
    line_width_pt: Optional[float] = None
    font_names: set[str] = field(default_factory=set)
    font_sizes: list[float] = field(default_factory=list)
    font_colors: set[str] = field(default_factory=set)
    bold_values: list[bool] = field(default_factory=list)
    italic_values: list[bool] = field(default_factory=list)
    alignments: list[str] = field(default_factory=list)
    runs: list[RunInfo] = field(default_factory=list)
    is_group_child: bool = False
    is_picture: bool = False
    is_connector: bool = False
    is_auto_shape: bool = False

    @property
    def right_cm(self) -> float:
        return self.left_cm + self.width_cm

    @property
    def bottom_cm(self) -> float:
        return self.top_cm + self.height_cm

    @property
    def center_x_cm(self) -> float:
        return self.left_cm + self.width_cm / 2

    @property
    def center_y_cm(self) -> float:
        return self.top_cm + self.height_cm / 2

    @property
    def area_cm2(self) -> float:
        return max(0.0, self.width_cm) * max(0.0, self.height_cm)


@dataclass
class XmlShapeInfo:
    tag: str
    name: str
    text: str
    left_cm: float
    top_cm: float
    width_cm: float
    height_cm: float
    geom: Optional[str] = None
    fill_color: Optional[str] = None
    line_color: Optional[str] = None
    line_width_pt: Optional[float] = None
    dashed: bool = False
    arrow_head: bool = False
    arrow_tail: bool = False
    flip_h: bool = False
    flip_v: bool = False
    group_depth: int = 0

    @property
    def right_cm(self) -> float:
        return self.left_cm + self.width_cm

    @property
    def bottom_cm(self) -> float:
        return self.top_cm + self.height_cm

    @property
    def center_x_cm(self) -> float:
        return self.left_cm + self.width_cm / 2

    @property
    def center_y_cm(self) -> float:
        return self.top_cm + self.height_cm / 2

    @property
    def length_cm(self) -> float:
        return math.hypot(self.width_cm, self.height_cm)

    @property
    def area_cm2(self) -> float:
        return max(0.0, self.width_cm) * max(0.0, self.height_cm)


@dataclass
class PptModel:
    path: Path
    valid_pptx: bool
    load_error: Optional[str]
    slide_count: int = 0
    slide_width_cm: float = 0.0
    slide_height_cm: float = 0.0
    shapes: list[ShapeInfo] = field(default_factory=list)
    xml_shapes: list[XmlShapeInfo] = field(default_factory=list)
    media_files: list[str] = field(default_factory=list)
    zip_entries: list[str] = field(default_factory=list)
    xml_counts: dict[str, int] = field(default_factory=dict)
    background_color: Optional[str] = None
    theme_colors: dict[str, str] = field(default_factory=dict)

    @property
    def all_text(self) -> str:
        texts = [shape.text for shape in self.shapes if shape.text.strip()]
        texts.extend(xml_shape.text for xml_shape in self.xml_shapes if xml_shape.text.strip())
        return normalize_space(" ".join(texts))

    @property
    def aspect_ratio(self) -> float:
        if self.slide_height_cm == 0:
            return 0.0
        return self.slide_width_cm / self.slide_height_cm


# Module-level cache of the currently-inspected theme's scheme->RGB map.
# color_is() consults this to resolve SCHEME:/THEME: references to real RGB
# so families are judged by the actual color, not by the theme role.
_ACTIVE_THEME_COLORS: dict[str, str] = {}


@dataclass
class CheckResult:
    rule_id: str
    label: str
    score: int
    passed: bool
    evidence: str


@dataclass
class Rule:
    rule_id: str
    label: str
    score: int
    check: Callable[[PptModel], CheckResult]


def normalize_space(value: str) -> str:
    return re.sub(r"\s+", " ", value or "").strip()


def normalize_for_match(value: str) -> str:
    return normalize_space(value).lower()


def cm_from_emu(value: int | float | None) -> float:
    return float(value or 0) / EMU_PER_CM


def pt_from_emu(value: int | float | None) -> Optional[float]:
    if value is None:
        return None
    return float(value) / EMU_PER_PT


def close(value: float, low: float, high: float, tol: float = CM_TOL) -> bool:
    return low - tol <= value <= high + tol


def in_spec_cm(value: float, low: float, high: float, tol: float = SPEC_CM_TOL) -> bool:
    """Return True when a measured dimension is inside the rubric range.

    Use this for explicit size/position ranges from the scoring rubric.  The
    tolerance is intentionally small and only absorbs EMU/cm conversion noise;
    the wider close() helper is for approximate layout proximity checks.
    """
    return low - tol <= value <= high + tol


def near(value: float, target: float, tol: float = CM_TOL) -> bool:
    return abs(value - target) <= tol


def safe_rgb(color_obj) -> Optional[str]:
    if color_obj is None:
        return None
    try:
        rgb = color_obj.rgb
        if rgb is not None:
            return str(rgb).upper()
    except Exception:
        pass
    try:
        theme = color_obj.theme_color
        if theme is not None:
            return f"THEME:{theme}"
    except Exception:
        pass
    return None


def shape_fill_color(shape) -> Optional[str]:
    try:
        return safe_rgb(shape.fill.fore_color)
    except Exception:
        return None


def shape_line_color(shape) -> Optional[str]:
    try:
        return safe_rgb(shape.line.color)
    except Exception:
        return None


def shape_line_width_pt(shape) -> Optional[float]:
    try:
        if shape.line.width is None:
            return None
        return pt_from_emu(shape.line.width)
    except Exception:
        return None


def effective_line_width_pt(width_pt: Optional[float]) -> float:
    """Return the rendered stroke width, substituting PowerPoint's default.

    OOXML omits <a:ln w="…"> for shapes that use the theme/master default
    outline, which PowerPoint draws at 1.0pt.  Rubric checks care about the
    visible stroke, so an unset width should be evaluated as 1.0pt rather
    than skipped or treated as "unknown".
    """
    return DEFAULT_LINE_WIDTH_PT if width_pt is None else width_pt


def local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1] if "}" in tag else tag


def find_first(element: ET.Element, path: str) -> Optional[ET.Element]:
    return element.find(path, NS)


def find_all(element: ET.Element, path: str) -> list[ET.Element]:
    return list(element.findall(path, NS))


def xml_attr_int(element: Optional[ET.Element], attr: str) -> int:
    if element is None:
        return 0
    try:
        return int(element.attrib.get(attr, "0"))
    except ValueError:
        return 0


def xml_text(element: ET.Element) -> str:
    return normalize_space("".join(t.text or "" for t in element.findall(".//a:t", NS)))


def xml_name(element: ET.Element) -> str:
    c_nv_pr = find_first(element, ".//p:cNvPr")
    return c_nv_pr.attrib.get("name", "") if c_nv_pr is not None else ""


def xml_geom(element: ET.Element) -> Optional[str]:
    geom = find_first(element, ".//a:prstGeom")
    return geom.attrib.get("prst") if geom is not None else None


def xml_color_from_parent(element: Optional[ET.Element]) -> Optional[str]:
    if element is None:
        return None
    srgb = find_first(element, ".//a:srgbClr")
    if srgb is not None and srgb.attrib.get("val"):
        return srgb.attrib["val"].upper()
    scheme = find_first(element, ".//a:schemeClr")
    if scheme is not None and scheme.attrib.get("val"):
        return f"SCHEME:{scheme.attrib['val']}"
    return None


def xml_fill_color(element: ET.Element) -> Optional[str]:
    sp_pr = find_first(element, ".//p:spPr")
    if sp_pr is None:
        sp_pr = find_first(element, ".//p:cxnSpPr")
    solid = find_first(sp_pr, ".//a:solidFill") if sp_pr is not None else None
    return xml_color_from_parent(solid)


def xml_background_color(root: ET.Element) -> Optional[str]:
    """Read the slide background fill color from <p:cSld>/<p:bg>.

    Returns the srgb/scheme color if the slide sets an explicit background,
    or None when the slide inherits its background (layout/master default),
    which PowerPoint renders as the default (white) background.
    """
    bg = find_first(root, ".//p:cSld/p:bg")
    if bg is None:
        return None
    return xml_color_from_parent(bg)


def xml_line_data(element: ET.Element) -> tuple[Optional[str], Optional[float], bool, bool, bool]:
    line = find_first(element, ".//a:ln")
    if line is None:
        return None, None, False, False, False
    color = xml_color_from_parent(line)
    width_pt = pt_from_emu(xml_attr_int(line, "w")) if "w" in line.attrib else None
    dash = find_first(line, "a:prstDash")
    dashed = dash is not None and dash.attrib.get("val", "solid") != "solid"
    head = find_first(line, "a:headEnd")
    tail = find_first(line, "a:tailEnd")
    arrow_head = head is not None and head.attrib.get("type", "none") != "none"
    arrow_tail = tail is not None and tail.attrib.get("type", "none") != "none"
    return color, width_pt, dashed, arrow_head, arrow_tail


def xml_position(element: ET.Element) -> tuple[float, float, float, float]:
    xfrm = find_first(element, ".//a:xfrm")
    off = find_first(xfrm, "a:off") if xfrm is not None else None
    ext = find_first(xfrm, "a:ext") if xfrm is not None else None
    return (
        cm_from_emu(xml_attr_int(off, "x")),
        cm_from_emu(xml_attr_int(off, "y")),
        cm_from_emu(xml_attr_int(ext, "cx")),
        cm_from_emu(xml_attr_int(ext, "cy")),
    )


def xml_flip(element: ET.Element) -> tuple[bool, bool]:
    """Return (flipH, flipV) of a shape's <a:xfrm>.

    PowerPoint draws a line from its bounding-box top-left to bottom-right by
    default.  flipV mirrors it vertically, which flips the visual slope: a box
    with height>0 that is flipV'd actually slants up-right instead of down-right.
    We need this to judge the rubric's up/down zig-zag directions correctly.
    """
    xfrm = find_first(element, ".//a:xfrm")
    if xfrm is None:
        return False, False
    return (
        xfrm.attrib.get("flipH", "0") in ("1", "true"),
        xfrm.attrib.get("flipV", "0") in ("1", "true"),
    )


def iter_xml_drawings(element: ET.Element, depth: int = 0) -> Iterable[tuple[ET.Element, int]]:
    for child in list(element):
        name = local_name(child.tag)
        if name in {"sp", "cxnSp", "pic", "grpSp"}:
            yield child, depth
            if name == "grpSp":
                yield from iter_xml_drawings(child, depth + 1)
        else:
            yield from iter_xml_drawings(child, depth)


def parse_theme_colors(theme_xml: bytes) -> dict[str, str]:
    """Parse ppt/theme/themeN.xml into a {scheme_name: 'RRGGBB'} map.

    We resolve srgbClr and sysClr@lastClr entries.  The resulting map lets
    color_is() judge scheme-colored fills by their actual RGB rather than by
    hardcoded scheme-role assumptions — a gray shape reported as accent3
    (A5A5A5) is judged gray because it is gray, not because "accent3" is on a
    hardcoded list.
    """
    result: dict[str, str] = {}
    try:
        root = ET.fromstring(theme_xml)
    except ET.ParseError:
        return result
    scheme = root.find(".//a:clrScheme", NS)
    if scheme is None:
        return result
    for child in list(scheme):
        name = local_name(child.tag)
        srgb = child.find(".//a:srgbClr", NS)
        if srgb is not None and srgb.attrib.get("val"):
            result[name.lower()] = srgb.attrib["val"].upper()
            continue
        sysc = child.find(".//a:sysClr", NS)
        if sysc is not None and sysc.attrib.get("lastClr"):
            result[name.lower()] = sysc.attrib["lastClr"].upper()
    return result


def parse_openxml(path: Path) -> tuple[list[str], list[str], list[XmlShapeInfo], dict[str, int], Optional[str], dict[str, str], Optional[str]]:
    if not zipfile.is_zipfile(path):
        return [], [], [], {}, None, {}, "not a valid pptx zip package"
    try:
        with zipfile.ZipFile(path) as archive:
            entries = archive.namelist()
            media_files = [name for name in entries if name.startswith("ppt/media/") and not name.endswith("/")]
            slide_xmls = sorted(
                name for name in entries if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
            )
            if not slide_xmls:
                return entries, media_files, [], {}, None, {}, "no slide XML found"
            slide_xml = archive.read(slide_xmls[0])
            theme_colors: dict[str, str] = {}
            theme_names = sorted(
                name for name in entries if re.fullmatch(r"ppt/theme/theme\d+\.xml", name)
            )
            if theme_names:
                try:
                    theme_colors = parse_theme_colors(archive.read(theme_names[0]))
                except Exception:
                    theme_colors = {}
    except Exception as exc:
        return [], [], [], {}, None, {}, str(exc)

    try:
        root = ET.fromstring(slide_xml)
    except ET.ParseError as exc:
        return entries, media_files, [], {}, None, {}, f"invalid slide XML: {exc}"

    background_color = xml_background_color(root)
    xml_shapes: list[XmlShapeInfo] = []
    counts = {"sp": 0, "cxnSp": 0, "pic": 0, "grpSp": 0}
    for element, depth in iter_xml_drawings(root):
        tag = local_name(element.tag)
        if tag in counts:
            counts[tag] += 1
        left, top, width, height = xml_position(element)
        line_color, line_width, dashed, arrow_head, arrow_tail = xml_line_data(element)
        flip_h, flip_v = xml_flip(element)
        xml_shapes.append(
            XmlShapeInfo(
                tag=tag,
                name=xml_name(element),
                text=xml_text(element),
                left_cm=left,
                top_cm=top,
                width_cm=width,
                height_cm=height,
                geom=xml_geom(element),
                fill_color=xml_fill_color(element),
                line_color=line_color,
                line_width_pt=line_width,
                dashed=dashed,
                arrow_head=arrow_head,
                arrow_tail=arrow_tail,
                flip_h=flip_h,
                flip_v=flip_v,
                group_depth=depth,
            )
        )
    return entries, media_files, xml_shapes, counts, background_color, theme_colors, None


def paragraph_alignment_name(paragraph) -> str:
    try:
        alignment = paragraph.alignment
        return str(alignment).lower() if alignment is not None else ""
    except Exception:
        return ""


def collect_text_style(shape) -> tuple[set[str], list[float], set[str], list[bool], list[bool], list[str], list[RunInfo]]:
    font_names: set[str] = set()
    font_sizes: list[float] = []
    font_colors: set[str] = set()
    bold_values: list[bool] = []
    italic_values: list[bool] = []
    alignments: list[str] = []
    runs: list[RunInfo] = []
    if not getattr(shape, "has_text_frame", False):
        return font_names, font_sizes, font_colors, bold_values, italic_values, alignments, runs
    for paragraph in shape.text_frame.paragraphs:
        alignments.append(paragraph_alignment_name(paragraph))
        # 段落级 defRPr（继承源）：当 run 属性缺失时回退取用，避免默认放行。
        para_font_name: Optional[str] = None
        para_size_pt: Optional[float] = None
        para_color: Optional[str] = None
        para_bold: Optional[bool] = None
        para_italic: Optional[bool] = None
        try:
            pf = paragraph.font
            try:
                if pf.name:
                    para_font_name = pf.name
            except Exception:
                pass
            try:
                if pf.size is not None:
                    para_size_pt = pt_from_emu(pf.size)
            except Exception:
                pass
            try:
                pc = safe_rgb(pf.color)
                if pc:
                    para_color = pc
            except Exception:
                pass
            try:
                if pf.bold is not None:
                    para_bold = bool(pf.bold)
            except Exception:
                pass
            try:
                if pf.italic is not None:
                    para_italic = bool(pf.italic)
            except Exception:
                pass
        except Exception:
            pass
        for run in paragraph.runs:
            font = run.font
            font_name = None
            size_pt = None
            color = None
            bold = None
            italic = None
            try:
                font_name = font.name
                if font_name:
                    font_names.add(font_name)
                elif para_font_name:
                    font_name = para_font_name
                    font_names.add(para_font_name)
            except Exception:
                pass
            try:
                if font.size is not None:
                    size_pt = pt_from_emu(font.size)
                    if size_pt is not None:
                        font_sizes.append(size_pt)
                elif para_size_pt is not None:
                    size_pt = para_size_pt
                    font_sizes.append(para_size_pt)
            except Exception:
                pass
            try:
                color = safe_rgb(font.color)
                if color:
                    font_colors.add(color)
                elif para_color:
                    color = para_color
                    font_colors.add(para_color)
            except Exception:
                pass
            try:
                bold = font.bold
                if bold is not None:
                    bold_values.append(bool(bold))
                elif para_bold is not None:
                    bold = para_bold
                    bold_values.append(para_bold)
            except Exception:
                pass
            try:
                italic = font.italic
                if italic is not None:
                    italic_values.append(bool(italic))
                elif para_italic is not None:
                    italic = para_italic
                    italic_values.append(para_italic)
            except Exception:
                pass
            runs.append(RunInfo(run.text, font_name, size_pt, color, bold, italic))
    return font_names, font_sizes, font_colors, bold_values, italic_values, alignments, runs


def shape_type_name(shape) -> str:
    try:
        return str(shape.shape_type).lower()
    except Exception:
        return type(shape).__name__


def extract_shape_info(shape, is_group_child: bool = False) -> list[ShapeInfo]:
    records: list[ShapeInfo] = []
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            records.extend(extract_shape_info(child, True))
    text = normalize_space(getattr(shape, "text", "") or "")
    font_names, font_sizes, font_colors, bold_values, italic_values, alignments, runs = collect_text_style(shape)
    shape_type = shape_type_name(shape)
    is_picture = False
    is_connector = False
    is_auto_shape = False
    try:
        if MSO_SHAPE_TYPE is not None:
            is_picture = shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            is_connector = shape.shape_type in {MSO_SHAPE_TYPE.LINE, MSO_SHAPE_TYPE.FREEFORM}
            is_auto_shape = shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    except Exception:
        pass
    records.append(
        ShapeInfo(
            name=getattr(shape, "name", ""),
            shape_type=shape_type,
            text=text,
            left_cm=cm_from_emu(getattr(shape, "left", 0)),
            top_cm=cm_from_emu(getattr(shape, "top", 0)),
            width_cm=cm_from_emu(getattr(shape, "width", 0)),
            height_cm=cm_from_emu(getattr(shape, "height", 0)),
            fill_color=shape_fill_color(shape),
            line_color=shape_line_color(shape),
            line_width_pt=shape_line_width_pt(shape),
            font_names=font_names,
            font_sizes=font_sizes,
            font_colors=font_colors,
            bold_values=bold_values,
            italic_values=italic_values,
            alignments=alignments,
            runs=runs,
            is_group_child=is_group_child,
            is_picture=is_picture,
            is_connector=is_connector,
            is_auto_shape=is_auto_shape,
        )
    )
    return records


def inspect_ppt(path: Path) -> PptModel:
    suffix = path.suffix.lower()
    if suffix != ".pptx":
        return PptModel(path, False, f"file extension {suffix or '(none)'} is not .pptx")
    if Presentation is None:
        return PptModel(path, False, "missing dependency: python-pptx; run python -m pip install -r requirements.txt")
    if not path.exists():
        return PptModel(path, False, "file does not exist")
    entries, media_files, xml_shapes, xml_counts, background_color, theme_colors, xml_error = parse_openxml(path)
    if xml_error:
        return PptModel(path, False, xml_error, media_files=media_files, zip_entries=entries, xml_shapes=xml_shapes, xml_counts=xml_counts, theme_colors=theme_colors)
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        return PptModel(path, False, f"python-pptx could not open file: {exc}", media_files=media_files, zip_entries=entries, xml_shapes=xml_shapes, xml_counts=xml_counts, theme_colors=theme_colors)

    shapes: list[ShapeInfo] = []
    if len(prs.slides) >= 1:
        for shape in prs.slides[0].shapes:
            shapes.extend(extract_shape_info(shape))
    _ACTIVE_THEME_COLORS.clear()
    _ACTIVE_THEME_COLORS.update(theme_colors)
    return PptModel(
        path=path,
        valid_pptx=True,
        load_error=None,
        slide_count=len(prs.slides),
        slide_width_cm=cm_from_emu(prs.slide_width),
        slide_height_cm=cm_from_emu(prs.slide_height),
        shapes=shapes,
        xml_shapes=xml_shapes,
        media_files=media_files,
        zip_entries=entries,
        xml_counts=xml_counts,
        background_color=background_color,
        theme_colors=theme_colors,
    )


def color_tuple(value: Optional[str]) -> Optional[tuple[int, int, int]]:
    if not value or not re.fullmatch(r"[0-9A-Fa-f]{6}", value):
        return None
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def color_distance(value: Optional[str], target: str) -> float:
    rgb = color_tuple(value)
    target_rgb = color_tuple(target)
    if rgb is None or target_rgb is None:
        return 999.0
    return math.sqrt(sum((a - b) ** 2 for a, b in zip(rgb, target_rgb)))


def _resolve_color(value: Optional[str]) -> Optional[str]:
    """Resolve a possibly scheme/theme-qualified color to a 6-hex RGB string.

    Returns None if value is None or unresolvable.  Direct RGB values pass
    through unchanged (uppercased).  SCHEME:xxx / THEME:xxx values are looked
    up in the active theme's color scheme so downstream family checks judge
    the actual pixel color, not the theme role.
    """
    if not value:
        return None
    value = value.upper()
    if value.startswith("SCHEME:") or value.startswith("THEME:"):
        key = value.split(":", 1)[1].lower()
        # PowerPoint's dk1/lt1/dk2/lt2 aliases show up as bg1/tx1/bg2/tx2 in
        # some documents; try both spellings.
        alias = {"bg1": "lt1", "tx1": "dk1", "bg2": "lt2", "tx2": "dk2"}.get(key, key)
        rgb = _ACTIVE_THEME_COLORS.get(key) or _ACTIVE_THEME_COLORS.get(alias)
        return rgb.upper() if rgb else None
    if re.fullmatch(r"[0-9A-F]{6}", value):
        return value
    return None


def color_is(value: Optional[str], family: str) -> bool:
    """Return True when `value` visually belongs to the named color family.

    Theme/scheme colors are resolved to their real RGB via the active theme
    before matching — the rubric only cares that a fill "is gray" (or blue,
    orange, etc.), not that it comes from a particular theme slot.  Only
    unresolved scheme references fall back to name-based heuristics.
    """
    if value is None:
        return False
    rgb_hex = _resolve_color(value)
    if rgb_hex is None:
        # Unresolvable scheme/theme reference — fall back to a coarse name
        # heuristic so we don't reject something we simply can't inspect.
        upper = value.upper()
        if upper.startswith("SCHEME:") or upper.startswith("THEME:"):
            scheme = upper.split(":", 1)[1].lower()
            if family == "white":
                return "lt1" in scheme or "bg1" in scheme or "lt2" in scheme
            if family == "black":
                return "dk1" in scheme or "tx1" in scheme
            return False
        return False
    rgb = color_tuple(rgb_hex)
    if rgb is None:
        return False
    r, g, b = rgb
    max_min = max(rgb) - min(rgb)
    if family == "white":
        return r >= 235 and g >= 235 and b >= 235
    if family == "black":
        return r <= 45 and g <= 45 and b <= 45
    if family == "gray":
        # Any neutral gray: channels close together, not near pure white/black.
        return max_min <= 25 and 40 <= (r + g + b) / 3 <= 220
    if family == "light_gray":
        return r >= 225 and g >= 225 and b >= 225 and max_min <= 18
    if family == "green":
        return g > r + 20 and g > b + 10 and g >= 85
    if family == "dark_green":
        return color_is(rgb_hex, "green") and sum(rgb) < 360
    if family == "blue":
        return b >= r and b >= g - 20 and b >= 90
    if family == "light_blue":
        return b >= 140 and g >= 130 and r >= 120 and b >= g >= r and max_min >= 12
    if family == "dark_teal":
        return g >= 45 and b >= 60 and r <= 85 and sum(rgb) < 380
    if family == "orange":
        return r >= 180 and 70 <= g <= 180 and b <= 90
    if family == "red":
        return r >= 170 and g <= 100 and b <= 100 and not color_is(rgb_hex, "orange")
    return False


def any_color_is(colors: Iterable[Optional[str]], family: str) -> bool:
    return any(color_is(color, family) for color in colors if color)


def text_shapes(model: PptModel, text: str, exact: bool = False) -> list[ShapeInfo]:
    needle = normalize_for_match(text)
    matches = []
    for shape in model.shapes:
        hay = normalize_for_match(shape.text)
        if exact and hay == needle:
            matches.append(shape)
        elif not exact and needle in hay:
            matches.append(shape)
    return matches


def xml_text_shapes(model: PptModel, text: str, exact: bool = False) -> list[XmlShapeInfo]:
    needle = normalize_for_match(text)
    matches = []
    for shape in model.xml_shapes:
        hay = normalize_for_match(shape.text)
        if exact and hay == needle:
            matches.append(shape)
        elif not exact and needle in hay:
            matches.append(shape)
    return matches


def all_text_contains(model: PptModel, text: str) -> bool:
    return normalize_for_match(text) in normalize_for_match(model.all_text)


def token_absent(model: PptModel, token: str) -> bool:
    return re.search(rf"\b{re.escape(token.lower())}\b", normalize_for_match(model.all_text)) is None


def average(values: list[float]) -> Optional[float]:
    return sum(values) / len(values) if values else None


def main_bbox(model: PptModel, include_small: bool = False) -> Optional[tuple[float, float, float, float]]:
    candidates = []
    for shape in model.shapes:
        if shape.width_cm <= 0 or shape.height_cm <= 0:
            continue
        if shape.is_picture:
            continue
        if not include_small and shape.area_cm2 < 0.08 and not shape.text:
            continue
        candidates.append((shape.left_cm, shape.top_cm, shape.right_cm, shape.bottom_cm))
    for shape in model.xml_shapes:
        if shape.tag == "pic" or shape.width_cm <= 0 or shape.height_cm <= 0:
            continue
        if not include_small and shape.area_cm2 < 0.08 and not shape.text:
            continue
        candidates.append((shape.left_cm, shape.top_cm, shape.right_cm, shape.bottom_cm))
    if not candidates:
        return None
    left = min(item[0] for item in candidates)
    top = min(item[1] for item in candidates)
    right = max(item[2] for item in candidates)
    bottom = max(item[3] for item in candidates)
    return left, top, right, bottom


def shapes_inside_slide(model: PptModel) -> bool:
    for shape in model.shapes:
        if shape.width_cm <= 0 or shape.height_cm <= 0:
            continue
        if shape.left_cm < -0.1 or shape.top_cm < -0.1:
            return False
        if shape.right_cm > model.slide_width_cm + 0.1 or shape.bottom_cm > model.slide_height_cm + 0.1:
            return False
    return True


def pass_result(rule_id: str, label: str, score: int, evidence: str) -> CheckResult:
    return CheckResult(rule_id, label, score, True, evidence)


def fail_result(rule_id: str, label: str, score: int, evidence: str) -> CheckResult:
    return CheckResult(rule_id, label, score, False, evidence)


def gate_valid_file(model: PptModel) -> CheckResult:
    if model.valid_pptx:
        return pass_result("D1-1", "交付文件为 .pptx 格式，文件可正常打开", 0, f"loaded successfully; zip entries={len(model.zip_entries)}")
    return fail_result("D1-1", "交付文件为 .pptx 格式，文件可正常打开", 0, model.load_error or "unknown load error")


def gate_not_single_image(model: PptModel) -> CheckResult:
    sp = model.xml_counts.get("sp", 0)
    connectors = model.xml_counts.get("cxnSp", 0)
    editable_lines = len(line_candidates(model))
    pics = model.xml_counts.get("pic", 0)
    text_count = sum(1 for shape in model.shapes if shape.text.strip())
    picture_area = sum(shape.area_cm2 for shape in model.xml_shapes if shape.tag == "pic")
    slide_area = max(model.slide_width_cm * model.slide_height_cm, 1)
    picture_ratio = picture_area / slide_area
    has_full_slide_picture = any(
        shape.tag == "pic"
        and shape.width_cm >= model.slide_width_cm * 0.90
        and shape.height_cm >= model.slide_height_cm * 0.90
        for shape in model.xml_shapes
    )
    has_editable_content = (sp + connectors + editable_lines) >= 5 or text_count >= 3
    ok = not has_full_slide_picture and picture_ratio < 0.85 and has_editable_content
    evidence = f"pics={pics}, picture_area={picture_ratio:.1%}, full_slide_picture={has_full_slide_picture}, editable_shapes={sp}, editable_lines={editable_lines}, text_shapes={text_count}"
    if ok:
        return pass_result("D1-3", "该幻灯片不是一整张图片", 0, evidence)
    return fail_result("D1-3", "该幻灯片不是一整张图片", 0, evidence)


def gate_check_replacement(model: PptModel) -> CheckResult:
    has_check = bool(text_shapes(model, "check")) or bool(xml_text_shapes(model, "check"))
    assessment_absent = token_absent(model, "assessment")
    if has_check and assessment_absent:
        return pass_result("D1-4", "第二部分标题为 check 且无 assessment", 0, "found title/text 'check'; no standalone 'assessment' token")
    return fail_result("D1-4", "第二部分标题为 check 且无 assessment", 0, f"has_check={has_check}, assessment_absent={assessment_absent}")


def gate_finish_body(model: PptModel) -> CheckResult:
    found = all_text_contains(model, FINISH_BODY)
    if found:
        return pass_result("D1-5", "finish 正文已替换为指定文本", 0, "exact normalized finish body found")
    return fail_result("D1-5", "finish 正文已替换为指定文本", 0, "specified finish body was not found after whitespace normalization")


def gate_filled_centered(model: PptModel) -> CheckResult:
    bbox = main_bbox(model)
    if bbox is None:
        return fail_result("D1-6", "整体构图铺满且居中无遮挡", 0, "no editable content bounding box")
    left, top, right, bottom = bbox
    width_ratio = (right - left) / model.slide_width_cm
    height_ratio = (bottom - top) / model.slide_height_cm
    center_dx = abs(((left + right) / 2) - model.slide_width_cm / 2) / model.slide_width_cm
    center_dy = abs(((top + bottom) / 2) - model.slide_height_cm / 2) / model.slide_height_cm
    ok = width_ratio >= 0.78 and height_ratio >= 0.68 and center_dx <= 0.08 and center_dy <= 0.10 and shapes_inside_slide(model)
    evidence = f"bbox={width_ratio:.1%}w x {height_ratio:.1%}h, center_dx={center_dx:.1%}, center_dy={center_dy:.1%}"
    if ok:
        return pass_result("D1-6", "整体构图铺满且居中无遮挡", 0, evidence)
    return fail_result("D1-6", "整体构图铺满且居中无遮挡", 0, evidence)


def stage_title_shape(model: PptModel, title: str) -> Optional[ShapeInfo]:
    matches = text_shapes(model, title, exact=True)
    if matches:
        return sorted(matches, key=lambda s: s.area_cm2, reverse=True)[0]
    matches = text_shapes(model, title)
    return sorted(matches, key=lambda s: s.area_cm2, reverse=True)[0] if matches else None


def stage_body_shape(model: PptModel, body: str) -> Optional[ShapeInfo]:
    matches = text_shapes(model, body)
    return sorted(matches, key=lambda s: len(s.text), reverse=True)[0] if matches else None


def number_shape(model: PptModel, number: str) -> Optional[ShapeInfo]:
    candidates = [shape for shape in model.shapes if normalize_space(shape.text) == number]
    if not candidates:
        return None
    return sorted(candidates, key=lambda s: (abs(s.width_cm - s.height_cm), -s.area_cm2))[0]


def nearby_shapes(model: PptModel, center_x: float, center_y: float, radius_cm: float) -> list[ShapeInfo]:
    found = []
    for shape in model.shapes:
        distance = math.hypot(shape.center_x_cm - center_x, shape.center_y_cm - center_y)
        if distance <= radius_cm:
            found.append(shape)
    return found


def nearby_xml_shapes(model: PptModel, center_x: float, center_y: float, radius_cm: float) -> list[XmlShapeInfo]:
    found = []
    for shape in model.xml_shapes:
        distance = math.hypot(shape.center_x_cm - center_x, shape.center_y_cm - center_y)
        if distance <= radius_cm:
            found.append(shape)
    return found


def shape_font_family_ok(shape: ShapeInfo, accepted: tuple[str, ...] = ("arial", "times", "liberation serif")) -> bool:
    if not shape.font_names:
        return True
    return any(any(accepted_name in name.lower() for accepted_name in accepted) for name in shape.font_names)


def font_size_in(shape: ShapeInfo, low: float, high: float) -> bool:
    sizes = [size for size in shape.font_sizes if size > 0]
    if not sizes:
        return True
    avg = average(sizes)
    return avg is not None and low - PT_TOL <= avg <= high + PT_TOL


def font_size_in_spec(shape: ShapeInfo, low: float, high: float, tol: float = 0.2) -> bool:
    """Check explicit rubric font-size ranges with only conversion tolerance."""
    sizes = [size for size in shape.font_sizes if size > 0]
    if not sizes:
        return True
    avg = average(sizes)
    return avg is not None and low - tol <= avg <= high + tol


def shape_has_text_color(shape: ShapeInfo, family: str) -> bool:
    if not shape.font_colors:
        return True
    return any_color_is(shape.font_colors, family)


def shape_bold_ok(shape: ShapeInfo, required: bool) -> bool:
    if not shape.bold_values:
        return True
    return any(shape.bold_values) if required else True


def shape_italic_ok(shape: ShapeInfo, required: bool) -> bool:
    if not shape.italic_values:
        return True
    return any(shape.italic_values) if required else True


def shape_centered_ok(shape: ShapeInfo) -> bool:
    if not shape.alignments:
        return True
    joined = " ".join(shape.alignments)
    return "center" in joined or "2" in joined or joined.strip() == ""


def find_title_bar(model: PptModel, title: str) -> Optional[ShapeInfo]:
    shape = stage_title_shape(model, title)
    if shape:
        return shape
    return None


def find_icon_candidates(model: PptModel, title_shape: Optional[ShapeInfo], number: str) -> list[ShapeInfo]:
    number_record = number_shape(model, number)
    anchor_x = title_shape.center_x_cm if title_shape else (number_record.center_x_cm if number_record else model.slide_width_cm / 2)
    anchor_y = (title_shape.top_cm - 1.8) if title_shape else (number_record.center_y_cm if number_record else model.slide_height_cm / 2)
    candidates = []
    for shape in model.shapes:
        if shape.text.strip():
            continue
        if 2.6 <= shape.width_cm <= 3.8 and 2.6 <= shape.height_cm <= 3.8:
            if abs(shape.center_x_cm - anchor_x) <= 1.2 and abs(shape.center_y_cm - anchor_y) <= 1.5:
                candidates.append(shape)
    if candidates:
        return candidates
    # Fallback to XML because PowerPoint often stores icon circles as background
    # shapes and the text/icon strokes as separate selectable objects.
    for shape in model.xml_shapes:
        if shape.tag == "sp" and shape.geom == "ellipse" and 2.6 <= shape.width_cm <= 3.8 and 2.6 <= shape.height_cm <= 3.8:
            if abs(shape.center_x_cm - anchor_x) <= 1.2 and abs(shape.center_y_cm - anchor_y) <= 1.5:
                candidates.append(
                    ShapeInfo(
                        name=shape.name,
                        shape_type="xml:sp",
                        text="",
                        left_cm=shape.left_cm,
                        top_cm=shape.top_cm,
                        width_cm=shape.width_cm,
                        height_cm=shape.height_cm,
                        fill_color=shape.fill_color,
                        line_color=shape.line_color,
                        line_width_pt=shape.line_width_pt,
                    )
                )
    return candidates


def xml_backing_shape(
    model: PptModel,
    text_shape: Optional[ShapeInfo],
    width_range: tuple[float, float],
    height_range: tuple[float, float],
    geoms: set[str],
    max_center_delta_cm: float = 0.45,
) -> Optional[XmlShapeInfo]:
    if text_shape is None:
        return None
    candidates = []
    for shape in model.xml_shapes:
        if shape.tag != "sp" or shape.text.strip():
            continue
        if shape.geom not in geoms:
            continue
        if not (close(shape.width_cm, width_range[0], width_range[1]) and close(shape.height_cm, height_range[0], height_range[1])):
            continue
        delta = math.hypot(shape.center_x_cm - text_shape.center_x_cm, shape.center_y_cm - text_shape.center_y_cm)
        if delta <= max_center_delta_cm:
            candidates.append((delta, shape))
    return sorted(candidates, key=lambda item: item[0])[0][1] if candidates else None


def text_effective_fill(model: PptModel, text_shape: Optional[ShapeInfo], width_range: tuple[float, float], height_range: tuple[float, float], geoms: set[str]) -> Optional[str]:
    if text_shape is None:
        return None
    return text_shape.fill_color or (xml_backing_shape(model, text_shape, width_range, height_range, geoms).fill_color if xml_backing_shape(model, text_shape, width_range, height_range, geoms) else None)


def number_backing_shape(model: PptModel, text_shape: Optional[ShapeInfo]) -> Optional[XmlShapeInfo]:
    return xml_backing_shape(model, text_shape, (0.9, 1.2), (0.9, 1.2), {"ellipse"}, 0.35)


def title_backing_shape(model: PptModel, text_shape: Optional[ShapeInfo]) -> Optional[XmlShapeInfo]:
    return xml_backing_shape(model, text_shape, (3.9, 4.5), (0.85, 1.15), {"roundRect"}, 0.45)


def body_backing_shape(model: PptModel, text_shape: Optional[ShapeInfo]) -> Optional[XmlShapeInfo]:
    return xml_backing_shape(model, text_shape, (5.1, 5.3), (4.3, 4.7), {"roundRect"}, 0.60)


def check_stage(model: PptModel, spec: dict) -> CheckResult:
    rule_id = f"D2-STAGE-{spec['num']}"
    label = f"第 {spec['num']} 阶段节点：{spec['title']}"
    score = 5
    # 标记为 strict 的阶段按 +5 细则逐点严格核验（细则未提及的不额外约束）。
    strict = bool(spec.get("strict"))
    failures = []
    num = number_shape(model, spec["num"])
    title = find_title_bar(model, spec["title"])
    body = stage_body_shape(model, spec["body"])
    if num is None:
        failures.append("missing number circle")
    else:
        num_bg = number_backing_shape(model, num)
        num_w = num_bg.width_cm if num_bg is not None else num.width_cm
        num_h = num_bg.height_cm if num_bg is not None else num.height_cm
        num_fill = num_bg.fill_color if num_bg is not None else num.fill_color
        # 细则：编号圆 宽高（spec['num_size'] 指定；未指定则沿用原默认 1-1.2cm）
        _ns = spec.get("num_size", (1.0, 1.2))
        num_lo = float(_ns[0]); num_hi = float(_ns[1])  # type: ignore[index,arg-type]
        if not (in_spec_cm(num_w, num_lo, num_hi) and in_spec_cm(num_h, num_lo, num_hi)):
            failures.append(f"number size={num_w:.2f}x{num_h:.2f}cm not in {num_lo}-{num_hi}cm")
        # 细则：深绿色填充
        if strict:
            if not color_is(num_fill, "dark_green"):
                failures.append(f"number fill not dark green: {num_fill}")
        elif not (color_is(num_fill, "dark_green") or color_is(num_fill, "green") or color_is(num_fill, "dark_teal")):
            failures.append(f"number fill={num_fill}")
        # 细则：白色编号（stage 1-4 要求；stage 5 细则未提"白色编号"，故不约束）
        if spec.get("num_white", True) and not shape_has_text_color(num, "white"):
            failures.append("number text is not white")
    if title is None:
        failures.append("missing title")
    else:
        title_bg = title_backing_shape(model, title)
        title_w = title_bg.width_cm if title_bg is not None else title.width_cm
        title_h = title_bg.height_cm if title_bg is not None else title.height_cm
        title_fill = title_bg.fill_color if title_bg is not None else title.fill_color
        # 细则：标题条 宽/高（spec['title_size'] 指定；未指定则沿用原默认 4-4.2cm × 0.9-1.1cm）
        t_w_lo, t_w_hi, t_h_lo, t_h_hi = spec.get("title_size", (4.0, 4.2, 0.9, 1.1))  # type: ignore[misc]
        t_w_lo = float(t_w_lo); t_w_hi = float(t_w_hi); t_h_lo = float(t_h_lo); t_h_hi = float(t_h_hi)
        if not (in_spec_cm(title_w, t_w_lo, t_w_hi) and in_spec_cm(title_h, t_h_lo, t_h_hi)):
            failures.append(f"title size={title_w:.2f}x{title_h:.2f}cm not in {t_w_lo}-{t_w_hi} x {t_h_lo}-{t_h_hi}cm")
        # 细则：深青色圆角标题条
        if strict:
            if not color_is(title_fill, "dark_teal"):
                failures.append(f"title fill not dark teal: {title_fill}")
            if title_bg is not None and title_bg.geom != "roundRect":
                failures.append(f"title not roundRect: {title_bg.geom}")
        elif not (color_is(title_fill, "dark_teal") or color_is(title_fill, "blue")):
            failures.append(f"title fill={title_fill}")
        # 细则：文本白色 Arial 或 新罗马(Times New Roman) 14-18磅 加粗
        if strict:
            if not shape_font_family_ok(title, ("arial", "times")):
                failures.append(f"title font not Arial/Times: {sorted(title.font_names)}")
            if not font_size_in_spec(title, 14, 18):
                failures.append(f"title font size not 14-18pt: {title.font_sizes}")
        else:
            if not shape_font_family_ok(title):
                failures.append(f"title font={sorted(title.font_names)}")
            if not font_size_in(title, 14, 18):
                failures.append(f"title font size={title.font_sizes}")
        if not shape_has_text_color(title, "white"):
            failures.append("title text not white")
        if not shape_bold_ok(title, True):
            failures.append("title not bold")
    if body is None:
        failures.append("missing body text")
    else:
        body_bg = body_backing_shape(model, body)
        body_w = body_bg.width_cm if body_bg is not None else body.width_cm
        body_h = body_bg.height_cm if body_bg is not None else body.height_cm
        body_fill = body_bg.fill_color if body_bg is not None else body.fill_color
        # 细则：说明框 宽/高（spec['body_size'] 指定；未指定则沿用原默认 5.1-5.3cm × 4.3-4.6cm）
        _bs = spec.get("body_size", (5.1, 5.3, 4.3, 4.6))
        b_w_lo = float(_bs[0]); b_w_hi = float(_bs[1]); b_h_lo = float(_bs[2]); b_h_hi = float(_bs[3])  # type: ignore[index,arg-type]
        if not (in_spec_cm(body_w, b_w_lo, b_w_hi) and in_spec_cm(body_h, b_h_lo, b_h_hi)):
            failures.append(f"body size={body_w:.2f}x{body_h:.2f}cm not in {b_w_lo}-{b_w_hi} x {b_h_lo}-{b_h_hi}cm")
        # 细则：圆角矩形说明框（strict 阶段要求 roundRect）
        if strict and body_bg is not None and body_bg.geom != "roundRect":
            failures.append(f"body not roundRect: {body_bg.geom}")
        # 细则：说明框填充色（stage 1 白色 / stage 2 浅蓝色，由 spec['body_fill'] 指定）
        if strict:
            want_fill = spec.get("body_fill")
            if want_fill and not color_is(body_fill, want_fill):
                failures.append(f"body fill not {want_fill}: {body_fill}")
        else:
            if not (color_is(body_fill, "light_blue") or color_is(body_fill, "white")):
                failures.append(f"body fill={body_fill}")
        # 细则：文本 Arial 或 新罗马 10-14磅（说明框正文内容由 stage_body_shape 匹配保证）
        if strict:
            if not shape_font_family_ok(body, ("arial", "times")):
                failures.append(f"body font not Arial/Times: {sorted(body.font_names)}")
            if not font_size_in_spec(body, 10, 14):
                failures.append(f"body font size not 10-14pt: {body.font_sizes}")
        else:
            if not shape_font_family_ok(body):
                failures.append(f"body font={sorted(body.font_names)}")
            if not font_size_in(body, 10, 14):
                failures.append(f"body font size={body.font_sizes}")
    icon_candidates = find_icon_candidates(model, title, spec["num"])
    if not icon_candidates:
        failures.append("missing editable icon circle near title")
    else:
        # 细则：浅蓝色圆形 宽高（spec['icon_size'] 指定；未指定则沿用原默认 3.1-3.4cm）
        _is = spec.get("icon_size", (3.1, 3.4))
        ic_lo = float(_is[0]); ic_hi = float(_is[1])  # type: ignore[index,arg-type]
        ic_mid = (ic_lo + ic_hi) / 2
        icon = sorted(icon_candidates, key=lambda s: abs(s.width_cm - ic_mid) + abs(s.height_cm - ic_mid))[0]
        if not (in_spec_cm(icon.width_cm, ic_lo, ic_hi) and in_spec_cm(icon.height_cm, ic_lo, ic_hi)):
            failures.append(f"icon size={icon.width_cm:.2f}x{icon.height_cm:.2f}cm not in {ic_lo}-{ic_hi}cm")
        if not color_is(icon.fill_color, "light_blue"):
            failures.append(f"icon fill={icon.fill_color}")
        # 细则：圆形内含图标 → 圆内存在可编辑图形部件
        local_xml = nearby_xml_shapes(model, icon.center_x_cm, icon.center_y_cm, 2.2)
        editable_inside = [x for x in local_xml if x.tag in {"sp", "cxnSp"} and x.area_cm2 < 8.0]
        if len(editable_inside) < 2:
            failures.append(f"icon editable parts={len(editable_inside)}")
        # 细则：圆形内含橙色感叹号小圆标（仅 stage 2）
        if spec.get("icon_orange_badge"):
            orange_badge = [
                x for x in local_xml
                if x.tag == "sp" and x.geom == "ellipse"
                and color_is(x.fill_color, "orange")
                and x.width_cm <= 1.2 and x.height_cm <= 1.2
            ]
            if not orange_badge:
                failures.append("orange exclamation badge circle not found in icon")
    # 细则：距上边线 3.5-3.8cm（仅 stage 1，spec['top_range'] 指定）
    if spec.get("top_range") and (title is not None or num is not None):
        top_low, top_high = spec["top_range"]
        num_bg = number_backing_shape(model, num) if num is not None else None
        stage_top = num_bg.top_cm if num_bg is not None else (num.top_cm if num is not None else title.top_cm)
        if strict:
            if not in_spec_cm(stage_top, top_low, top_high):
                failures.append(f"stage top={stage_top:.2f}cm not in {top_low}-{top_high}cm")
        elif not close(stage_top, top_low, top_high, 0.15):
            failures.append(f"stage top={stage_top:.2f}cm, expected near {top_low}-{top_high}cm")
    # 细则：位于指定阶段右侧（stage 2 位于第 1 阶段 start 右侧）
    right_of = spec.get("right_of_title")
    if right_of and title is not None:
        left_stage = stage_title_shape(model, right_of)
        if left_stage is not None and not (title.center_x_cm > left_stage.center_x_cm):
            failures.append(f"stage not to the right of '{right_of}' (cx={title.center_x_cm:.2f} <= {left_stage.center_x_cm:.2f})")
    # 细则：位于页面中部（stage 3）——节点水平中心接近页面水平中线
    if spec.get("center_of_page") and title is not None:
        page_center_x = model.slide_width_cm / 2
        if not near(title.center_x_cm, page_center_x, 1.0):
            failures.append(f"stage not at page center (cx={title.center_x_cm:.2f} vs {page_center_x:.2f})")
    # 细则：位于页面右侧（stage 5）——节点水平中心位于页面右半区
    if spec.get("right_side_of_page") and title is not None:
        if not (title.center_x_cm > model.slide_width_cm * 0.6):
            failures.append(f"stage not on right side of page (cx={title.center_x_cm:.2f}, page_w={model.slide_width_cm:.2f})")
    if failures:
        return fail_result(rule_id, label, score, "; ".join(failures))
    return pass_result(rule_id, label, score, f"found number {spec['num']}, title '{spec['title']}', body, icon and styles")


def line_candidates(model: PptModel) -> list[XmlShapeInfo]:
    return [shape for shape in model.xml_shapes if shape.tag == "cxnSp" or (shape.line_width_pt is not None and shape.width_cm + shape.height_cm > 0.2)]


def is_straight_line_geom(shape: XmlShapeInfo) -> bool:
    """Return True when a shape is a single-stroke straight line/arrow.

    Accepts both routes PowerPoint uses for "一体箭头":
      • autoshape ``line`` (``<p:sp>`` with ``prstGeom prst="line"``);
      • the straight-connector family (``<p:cxnSp>`` with ``straightConnector1``
        or plain ``line`` geometry).
    Either can carry ``headEnd``/``tailEnd`` arrowheads as one integrated
    object, so both should be eligible everywhere an "arrow" is checked.
    """
    if shape.geom == "line" or shape.geom == "straightConnector1":
        return True
    # Straight connectors sometimes drop geom in the XML tree; treat any
    # cxnSp with near-zero thickness (a rendered straight stroke) as one too.
    if shape.tag == "cxnSp" and (shape.geom is None or "straight" in (shape.geom or "").lower()):
        return True
    return False


def action_arrow_candidates(model: PptModel) -> list[XmlShapeInfo]:
    candidates = []
    for shape in line_candidates(model):
        horizontal = shape.width_cm >= 2.0 and abs(shape.height_cm) <= 0.6
        dark_line = color_is(shape.line_color, "black") or color_is(shape.line_color, "dark_teal") or shape.line_color is None
        # Unset line widths render at the default 1.0pt; judge that instead of skipping.
        width = effective_line_width_pt(shape.line_width_pt)
        width_ok = 0.8 <= width <= 1.8
        arrow = shape.arrow_head or shape.arrow_tail
        upper_flow = 4.0 <= shape.center_y_cm <= 8.0
        if horizontal and dark_line and width_ok and arrow and not shape.dashed and upper_flow:
            candidates.append(shape)
    return sorted(candidates, key=lambda s: (s.top_cm, s.left_cm))


def emotion_line_candidates(model: PptModel) -> list[XmlShapeInfo]:
    candidates = []
    for shape in line_candidates(model):
        orange = color_is(shape.line_color, "orange")
        width = effective_line_width_pt(shape.line_width_pt)
        width_ok = 0.8 <= width <= 2.3
        lower = shape.center_y_cm >= model.slide_height_cm * 0.45
        short_dash_piece = is_straight_line_geom(shape) and 0.05 <= shape.length_cm <= 0.45
        if orange and lower and width_ok and (shape.dashed or short_dash_piece):
            candidates.append(shape)
    return sorted(candidates, key=lambda s: s.left_cm)


def score_legend(model: PptModel) -> CheckResult:
    # 细则 (+3)：第1页左上图例框
    #   ① 位于页面左上角
    #   ② 距上边线 0.2-1cm
    #   ③ 距左边线 0.4-1cm
    #   ④ 宽度 5-6.5cm
    #   ⑤ 高度 2-2.6cm
    #   ⑥ 白色填充
    #   ⑦ 深蓝色圆角矩形边框
    #   ⑧ 长度 2-2.6cm 黑色实线箭头
    #   ⑨ 文本 "action"
    #   ⑩ 长度 2-2.6cm 橙色虚线箭头
    #   ⑪ 文本 "emotion"
    #   ⑫ 文本字体 Arial 或 Times New Roman(新罗马)，14-18 磅
    # 仅检查细则列出的点，不额外约束；判定基于 PowerPoint 实际存储的形状与线条。
    failures: list[str] = []

    action = text_shapes(model, "action")
    emotion = text_shapes(model, "emotion")

    def in_box_spec(shape: XmlShapeInfo) -> bool:
        # ②③④⑤ 位置与尺寸区间（左上角由左/上边距区间共同保证 → ①）
        return (
            in_spec_cm(shape.top_cm, 0.2, 1.0)
            and in_spec_cm(shape.left_cm, 0.4, 1.0)
            and in_spec_cm(shape.width_cm, 5.0, 6.5)
            and in_spec_cm(shape.height_cm, 2.0, 2.6)
        )

    # ⑦ 图例框必须是圆角矩形（roundRect）。取底层 XML 形状（含 geom / 填充 / 边框）。
    legend_boxes = [
        shape for shape in model.xml_shapes
        if shape.tag == "sp" and shape.geom == "roundRect" and in_box_spec(shape)
    ]
    box = legend_boxes[0] if legend_boxes else None
    if box is None:
        failures.append("legend roundRect box (pos/size) not found")

    if box is not None:
        # ⑥ 白色填充
        if not color_is(box.fill_color, "white"):
            failures.append(f"box fill not white: {box.fill_color}")
        # ⑦ 深蓝色边框
        def is_dark_blue(color: Optional[str]) -> bool:
            rgb = color_tuple(color)
            if rgb is None:
                return color_is(color, "blue")
            r, g, b = rgb
            return b >= r and b >= g and b > 40 and sum(rgb) < 420
        if not is_dark_blue(box.line_color):
            failures.append(f"box border not dark blue: {box.line_color}")

    # 箭头必须位于图例框范围内（细则：图例框"包含"两种箭头）。
    def inside_legend(line: XmlShapeInfo) -> bool:
        if box is None:
            return False
        return (
            box.left_cm - 0.2 <= line.center_x_cm <= box.right_cm + 0.2
            and box.top_cm - 0.2 <= line.center_y_cm <= box.bottom_cm + 0.2
        )

    legend_lines = [line for line in line_candidates(model) if is_straight_line_geom(line) and inside_legend(line)]

    # ⑧ 长度 2-2.6cm 黑色实线箭头
    black_arrows = [
        line for line in legend_lines
        if (line.arrow_head or line.arrow_tail)
        and not line.dashed
        and color_is(line.line_color, "black")
        and in_spec_cm(line.length_cm, 2.0, 2.6)
    ]
    if not black_arrows:
        failures.append("black solid arrow (len 2-2.6cm) not found")

    # ⑩ 长度 2-2.6cm 橙色虚线箭头
    orange_arrows = [
        line for line in legend_lines
        if (line.arrow_head or line.arrow_tail)
        and line.dashed
        and color_is(line.line_color, "orange")
        and in_spec_cm(line.length_cm, 2.0, 2.6)
    ]
    if not orange_arrows:
        failures.append("orange dashed arrow (len 2-2.6cm) not found")

    # ⑨⑪ 文本 action / emotion
    if not action:
        failures.append("text 'action' not found")
    if not emotion:
        failures.append("text 'emotion' not found")

    # ⑫ 文本字体 Arial 或 Times New Roman，14-18 磅
    for shape in action + emotion:
        if not shape_font_family_ok(shape, ("arial", "times")):
            failures.append(f"legend text font not Arial/Times: {sorted(shape.font_names)}")
        if not font_size_in_spec(shape, 14, 18):
            failures.append(f"legend text size not 14-18pt: {shape.font_sizes}")

    evidence = (
        f"box={box is not None}, "
        f"fill={box.fill_color if box else None}, border={box.line_color if box else None}, "
        f"black_arrows={len(black_arrows)}, orange_dashed_arrows={len(orange_arrows)}, "
        f"action_text={len(action)}, emotion_text={len(emotion)}"
    )
    if failures:
        return fail_result("D2-02", "左上图例框", 3, "; ".join(failures) + f" | {evidence}")
    return pass_result("D2-02", "左上图例框", 3, evidence)


def score_number_style(model: PptModel) -> CheckResult:
    # 细则 (+1)：编号 "1""2""3""4""5" 为 Arial 或 新罗马(Times New Roman)
    #   18-22磅、白色、加粗。仅检查这些点，不额外约束。
    # 说明：样式解析（collect_text_style）已支持段落级 defRPr 继承——run 属性
    #   缺失时会自动回退到段落 font 的对应属性并写入统一字段。仍无可读证据
    #   时不再强制判失败，交由 helper 保持"未显式设置 → 继承默认"的宽容判定。
    failures = []
    for number in ["1", "2", "3", "4", "5"]:
        shape = number_shape(model, number)
        if shape is None:
            failures.append(f"missing {number}")
            continue
        # 字体：Arial 或 新罗马（未读取到字体信息时视为继承母版默认，宽容通过）
        if not shape_font_family_ok(shape, ("arial", "times")):
            failures.append(f"{number} font not Arial/Times: {sorted(shape.font_names)}")
        # 字号：18-22磅（未读取到字号时视为继承母版默认，宽容通过）
        if not font_size_in_spec(shape, 18, 22):
            failures.append(f"{number} size not 18-22pt: {shape.font_sizes}")
        # 颜色：白色（未读取到字体颜色时视为继承母版默认，宽容通过）
        if not shape_has_text_color(shape, "white"):
            failures.append(f"{number} not white: {sorted(shape.font_colors)}")
        # 加粗（未读取到加粗设置时视为继承母版默认，宽容通过）
        if not shape_bold_ok(shape, True):
            failures.append(f"{number} not bold")
    if failures:
        return fail_result("D2-08", "编号样式", 1, "; ".join(failures))
    return pass_result("D2-08", "编号样式", 1, "numbers 1-5 are Arial/Times, 18-22pt, white, bold")


def score_action_arrows(model: PptModel) -> CheckResult:
    # 细则 (+1)：第1至第5阶段之间设置 4 条黑色水平单箭头，方向均从左向右，
    #   线宽 1.0-3磅，长度 2-3cm，箭头连接相邻阶段图标区域。
    #   仅检查细则列出的点，不额外约束；判定基于 PowerPoint 实际存储的线条。
    # 视觉方向计算：
    #   OOXML 中线段起点位于形状 bbox 左上 (left, top)、终点位于右下 (right, bottom)。
    #   flipH 会左右翻转：起点变到右上、终点变到左下（视觉方向水平反向）。
    #   flipV 会上下翻转：不影响水平方向。
    #   arrow_head 附着于形状起点，arrow_tail 附着于形状终点。
    #   由此，箭头视觉终点的 x 坐标计算方式：
    #     ┌ 无 flipH：起点在左、终点在右 → tail 在右、head 在左
    #     └ 有 flipH：起点在右、终点在左 → tail 在左、head 在右
    #   视觉"从左向右"即"箭头尖端位于右端"，等价条件：
    #     (arrow_tail 且 !flipH) 或 (arrow_head 且 flipH)

    # 各阶段图标圆心（按从左到右排序），用于校验"连接相邻阶段图标区域"。
    icon_centers: list[tuple[float, float]] = []
    for spec in STAGE_SPECS:
        title = stage_title_shape(model, str(spec["title"]))
        icons = find_icon_candidates(model, title, str(spec["num"]))
        if icons:
            icon = sorted(icons, key=lambda s: abs(s.width_cm - 3.2) + abs(s.height_cm - 3.2))[0]
            icon_centers.append((icon.center_x_cm, icon.center_y_cm))
    icon_centers.sort(key=lambda c: c[0])

    # 逐条筛选满足全部细则要求的黑色水平单箭头。
    valid_arrows: list[XmlShapeInfo] = []
    for shape in line_candidates(model):
        # 一体箭头：autoshape line 或直线连接符 straightConnector1 均视为等价
        if not is_straight_line_geom(shape):
            continue
        # 水平
        if abs(shape.height_cm) > 0.15:
            continue
        # 黑色
        if not color_is(shape.line_color, "black"):
            continue
        # 单箭头（只有一端有箭头）
        single_arrow = shape.arrow_head != shape.arrow_tail
        if not single_arrow:
            continue
        # 方向从左向右（视觉）：基于起点/终点/flipH 与箭头端综合判定，避免仅凭
        # tailEnd 导致的反向误判（形状被水平翻转时 tailEnd 反而指向左侧）。
        if shape.flip_h:
            # flipH：起点在右、终点在左 → arrow_head(起点端) 指向右侧
            visual_left_to_right = shape.arrow_head and not shape.arrow_tail
        else:
            # 无 flipH：起点在左、终点在右 → arrow_tail(终点端) 指向右侧
            visual_left_to_right = shape.arrow_tail and not shape.arrow_head
        if not visual_left_to_right:
            continue
        # 线宽 1.0-3磅（未显式设置线宽时按 PowerPoint 默认 1.0pt 判定）
        width = effective_line_width_pt(shape.line_width_pt)
        if not (1.0 - PT_TOL <= width <= 3.0 + PT_TOL):
            continue
        # 长度 2-3cm
        if not in_spec_cm(shape.length_cm, 2.0, 3.0):
            continue
        valid_arrows.append(shape)

    # 箭头连接相邻阶段图标区域：箭头位于相邻两图标之间（水平夹在两圆心间、
    # 纵向与图标同高），且 4 个相邻间隔各有一条。
    connected_gaps = 0
    if len(icon_centers) == 5:
        for i in range(4):
            lx, ly = icon_centers[i]
            rx, _ = icon_centers[i + 1]
            for arrow in valid_arrows:
                if lx < arrow.center_x_cm < rx and abs(arrow.center_y_cm - ly) <= 1.5:
                    connected_gaps += 1
                    break

    ok = len(valid_arrows) == 4 and connected_gaps == 4
    evidence = (
        f"valid_arrows={len(valid_arrows)} (need 4 black horizontal single arrows, "
        f"L->R, 1.0-3pt, 2-3cm), connected_adjacent_icon_gaps={connected_gaps}/4"
    )
    return pass_result("D2-09", "action 主箭头", 1, evidence) if ok else fail_result("D2-09", "action 主箭头", 1, evidence)


def score_emotion_line(model: PptModel) -> CheckResult:
    # 细则 (+3)：页面下方橙色虚线折线箭头（可编辑一体折线）
    #   6 段斜线从左到右依次为：
    #     ① 2.5-3.5cm 斜向右下方
    #     ② 8-9.5cm  斜向右上方
    #     ③ 6-7cm    斜向右下方
    #     ④ 4.5-6cm  斜向右上方
    #     ⑤ 2-3cm    斜向右下方
    #     ⑥ 6-7cm    斜向右上方（箭头结束）
    #   线宽 1.0-3磅，虚线样式清晰。仅检查细则列出的点，不额外约束。
    lines = emotion_line_candidates(model)
    if not lines:
        return fail_result("D2-10", "emotion 橙色虚线折线", 3, "no orange dashed lower-page line found")

    # The rubric describes one editable dashed zig-zag made of six diagonal
    # segments with explicit lengths and line width.  Do not accept many tiny
    # solid line pieces that only visually imitate a dashed line.
    long_segments = [
        line for line in lines
        if is_straight_line_geom(line) and line.length_cm > 1.0
    ]
    sorted_lines = sorted(long_segments, key=lambda s: s.left_cm)[:6]
    expected_ranges = [(2.5, 3.5), (8.0, 9.5), (6.0, 7.0), (4.5, 6.0), (2.0, 3.0), (6.0, 7.0)]
    lengths = [line.length_cm for line in sorted_lines]
    length_ok = len(sorted_lines) == 6 and all(
        in_spec_cm(length, low, high, 0.05)
        for length, (low, high) in zip(lengths, expected_ranges)
    )
    directions = [1 if line.height_cm > 0 else -1 for line in sorted_lines]
    direction_ok = directions == [1, -1, 1, -1, 1, -1]
    width_ok = len(sorted_lines) == 6 and all(
        1.0 <= effective_line_width_pt(line.line_width_pt) <= 3.0
        for line in sorted_lines
    )
    dashed_ok = len(sorted_lines) == 6 and all(line.dashed for line in sorted_lines)
    arrow_end = bool(sorted_lines) and (sorted_lines[-1].arrow_head or sorted_lines[-1].arrow_tail)
    ok = length_ok and direction_ok and width_ok and dashed_ok and arrow_end
    short_piece_count = sum(1 for line in lines if line.length_cm <= 0.45)
    evidence = (
        f"segments={len(long_segments)}, lengths={[round(v, 2) for v in lengths]} "
        f"(need 2.5-3.5/8-9.5/6-7/4.5-6/2-3/6-7cm), "
        f"length_ok={length_ok}, directions={directions}, direction_ok={direction_ok}, "
        f"width_ok={width_ok} (need 1.0-3pt), dashed_ok={dashed_ok}, arrow_end={arrow_end}, "
        f"short_pieces={short_piece_count}"
    )
    return pass_result("D2-10", "emotion 橙色虚线折线", 3, evidence) if ok else fail_result("D2-10", "emotion 橙色虚线折线", 3, evidence)


def score_emotion_labels(model: PptModel) -> CheckResult:
    # 细则 (+1)：第1页情绪标签
    #   ① 底部依次出现 "concerned" "tense and focused" "hopeful" "relieved"（从左到右顺序）
    #   ② 字体为 Times New Roman 或等效西文字体
    #   ③ 字号 12-15 磅
    #   ④ 橙色
    #   ⑤ 倾斜（italic）
    # 仅检查细则列出的点，不额外约束；判定基于 PowerPoint 实际存储的文本样式与位置。
    failures: list[str] = []
    label_shapes: list[ShapeInfo] = []
    for word in EMOTION_WORDS:
        matches = text_shapes(model, word)
        if not matches:
            failures.append(f"missing {word}")
            continue
        shape = matches[0]
        label_shapes.append(shape)
        # ② 字体 Times New Roman 或等效西文字体
        if not shape_font_family_ok(shape, ("times", "liberation serif", "cambria", "georgia")):
            failures.append(f"{word} font={sorted(shape.font_names)}")
        # ③ 字号 12-15 磅
        if not font_size_in_spec(shape, 12, 15):
            failures.append(f"{word} size not 12-15pt: {shape.font_sizes}")
        # ④ 橙色
        if not shape_has_text_color(shape, "orange"):
            failures.append(f"{word} not orange")
        # ⑤ 倾斜
        if not shape_italic_ok(shape, True):
            failures.append(f"{word} not italic")
        # ① 底部
        if shape.center_y_cm < model.slide_height_cm * 0.45:
            failures.append(f"{word} not at bottom y={shape.center_y_cm:.2f}cm")
    # ① 依次出现（从左到右顺序）
    if len(label_shapes) == len(EMOTION_WORDS):
        centers = [shape.center_x_cm for shape in label_shapes]
        if centers != sorted(centers):
            failures.append(f"labels not left-to-right order={centers}")
    if failures:
        return fail_result("D2-11", "情绪标签", 1, "; ".join(failures))
    return pass_result("D2-11", "情绪标签", 1, "all four emotion labels found at bottom with Times/西文 12-15pt, orange, italic")


def score_theme_text(model: PptModel) -> CheckResult:
    # 细则 (+1)：第1页底部主题文字
    #   ① 页面底部中央出现文本 "theme: managing a website outage before a product launch"
    #   ② 字体 Times New Roman 或等效西文字体
    #   ③ 字号 15-19 磅
    #   ④ 蓝色
    #   ⑤ 倾斜（italic）
    # 仅检查细则列出的点，不额外约束；判定基于 PowerPoint 实际存储的文本、样式与位置。
    matches = text_shapes(model, THEME_TEXT)
    if not matches:
        return fail_result("D2-12", "底部主题文字", 1, "theme text not found")
    shape = matches[0]
    failures = []
    # ① 页面底部
    if shape.center_y_cm < model.slide_height_cm * 0.84:
        failures.append(f"not at bottom y={shape.center_y_cm:.2f}cm")
    # ① 页面中央（水平方向文本框居中于页面）
    if abs(shape.center_x_cm - model.slide_width_cm / 2) > 1.5:
        failures.append(f"not centered x={shape.center_x_cm:.2f}cm")
    # ② 字体 Times New Roman 或等效西文字体
    if not shape_font_family_ok(shape, ("times", "liberation serif", "cambria", "georgia")):
        failures.append(f"font={sorted(shape.font_names)}")
    # ③ 字号 15-19 磅
    if not font_size_in_spec(shape, 15, 19):
        failures.append(f"size not 15-19pt: {shape.font_sizes}")
    # ④ 蓝色
    if not shape_has_text_color(shape, "blue"):
        failures.append("not blue")
    # ⑤ 倾斜
    if not shape_italic_ok(shape, True):
        failures.append("not italic")
    if failures:
        return fail_result("D2-12", "底部主题文字", 1, "; ".join(failures))
    return pass_result("D2-12", "底部主题文字", 1, "theme text at bottom center with Times/西文 15-19pt, blue, italic")


def score_background(model: PptModel) -> CheckResult:
    # 细则 (+1)：第1页背景
    #   ① 页面背景为白色或极浅灰白色
    # 仅检查细则列出的点，不额外约束；判定基于 PowerPoint 实际存储的背景填充。
    failures: list[str] = []

    # ① 背景为白色或极浅灰白色。未显式设置背景时继承母版默认（PowerPoint 渲染为白色），视为通过。
    bg = model.background_color
    if bg is not None and not (color_is(bg, "white") or color_is(bg, "light_gray")):
        failures.append(f"background not white/near-white ({bg})")

    evidence = f"background={bg or 'inherited(default white)'}"
    if failures:
        return fail_result("D2-14", "背景干净", 1, "; ".join(failures))
    return pass_result("D2-14", "背景干净", 1, evidence)


def build_gate_rules() -> list[Rule]:
    return [
        Rule("D1-1", "交付文件为 .pptx 格式，文件可正常打开", 0, gate_valid_file),
        Rule("D1-3", "该幻灯片不是一整张图片", 0, gate_not_single_image),
    ]


def make_stage_rule(spec: dict) -> Rule:
    return Rule(f"D2-STAGE-{spec['num']}", f"第 {spec['num']} 阶段节点：{spec['title']}", 5, lambda model, spec=spec: check_stage(model, spec))


def build_scoring_rules() -> list[Rule]:
    rules = [
        Rule("D2-02", "左上图例框", 3, score_legend),
    ]
    rules.extend(make_stage_rule(spec) for spec in STAGE_SPECS)
    rules.extend(
        [
            Rule("D2-08", "编号样式", 1, score_number_style),
            Rule("D2-09", "action 主箭头", 1, score_action_arrows),
            Rule("D2-10", "emotion 橙色虚线折线", 3, score_emotion_line),
            Rule("D2-11", "情绪标签", 1, score_emotion_labels),
            Rule("D2-12", "底部主题文字", 1, score_theme_text),
            Rule("D2-14", "背景干净", 1, score_background),
        ]
    )
    return rules


def run_rules(model: PptModel, rules: list[Rule]) -> list[CheckResult]:
    results = []
    for rule in rules:
        try:
            results.append(rule.check(model))
        except Exception as exc:
            results.append(fail_result(rule.rule_id, rule.label, rule.score, f"check raised {type(exc).__name__}: {exc}"))
    return results


def _run_evaluation(path: Path) -> tuple[PptModel, list[CheckResult], list[CheckResult], int, bool]:
    model = inspect_ppt(path)
    gate_results = run_rules(model, build_gate_rules())
    passed_gate = all(result.passed for result in gate_results)
    if not passed_gate:
        return model, gate_results, [], 0, False
    scoring_results = run_rules(model, build_scoring_rules())
    total = sum(result.score for result in scoring_results if result.passed)
    return model, gate_results, scoring_results, total, True


def _locate_pptx(dir_path: Path) -> Optional[Path]:
    """Locate the PPTX document to evaluate inside `dir_path`.

    Prefers the well-known filename shipped alongside the script; falls back
    to a single .pptx in the directory when the well-known name is absent,
    so the script keeps working if the document is renamed.
    """
    preferred = dir_path / "website_outage_editable_flow.pptx"
    if preferred.is_file():
        return preferred
    candidates = sorted(
        p for p in dir_path.iterdir()
        if p.is_file() and p.suffix.lower() == ".pptx"
    )
    return candidates[0] if candidates else None


# 维度二评分项定义：与 build_scoring_rules() 顺序一致，用于统一输出结构中的 dim2_items。
_DIM2_ITEMS: list[tuple[str, str, int]] = [
    ("D2-02", "左上图例框", 3),
    *((f"D2-STAGE-{spec['num']}", f"第 {spec['num']} 阶段节点：{spec['title']}", 5) for spec in STAGE_SPECS),
    ("D2-08", "编号样式", 1),
    ("D2-09", "action 主箭头", 1),
    ("D2-10", "emotion 橙色虚线折线", 3),
    ("D2-11", "情绪标签", 1),
    ("D2-12", "底部主题文字", 1),
    ("D2-14", "背景干净", 1),
]

_MAX_SCORE = sum(max_delta for _, _, max_delta in _DIM2_ITEMS)


def evaluate(dir_path: str) -> dict:
    """Evaluate the PPT located in `dir_path` and return a structured report.

    `dir_path` is the directory where this script resides; the script locates
    and opens the target document inside that directory itself.  Returns a
    dict following the batch-runner contract (see §2.2 of the interface spec).
    """
    result: dict = {
        "id": "055",
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": _MAX_SCORE,
    }
    try:
        directory = Path(dir_path)
        if not directory.is_dir():
            result["status"] = "error"
            result["error"] = f"dir_path is not a directory: {dir_path}"
            return result
        target = _locate_pptx(directory)
        if target is None:
            result["status"] = "error"
            result["error"] = f"no .pptx file found in {dir_path}"
            return result
        result["file_name"] = target.name

        model, gate_results, scoring_results, total, passed_gate = _run_evaluation(target)
        result["dim1_pass"] = passed_gate
        if not passed_gate:
            reasons = [
                f"{r.rule_id} {r.label}: {r.evidence}"
                for r in gate_results if not r.passed
            ]
            result["dim1_reason"] = "; ".join(reasons)
            return result

        scoring_by_id = {r.rule_id: r for r in scoring_results}
        items = []
        for rule_id, label, max_delta in _DIM2_ITEMS:
            check = scoring_by_id.get(rule_id)
            hit = bool(check and check.passed)
            items.append({
                "rule": label,
                "max_delta": max_delta,
                "delta": max_delta if hit else 0,
                "hit": hit,
                "detail": "",
            })
        result["dim2_items"] = items
        result["total_score"] = total
        return result
    except Exception as exc:
        result["status"] = "error"
        result["error"] = f"{type(exc).__name__}: {exc}"
        return result


if __name__ == "__main__":
    _arg = sys.argv[1] if len(sys.argv) > 1 else str(Path(__file__).resolve().parent)
    print(json.dumps(evaluate(_arg), ensure_ascii=False, indent=2))
