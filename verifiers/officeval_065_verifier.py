# -*- coding: utf-8 -*-
"""
对 "口算教学课件_第10页双转盘动画版.pptx" 的自动评估脚本

评分逻辑
========
维度 1（可用与可修改性）：
    - 文件为 .pptx
    - 能正常被解析（即能被打开的代理）
    - 第 10 页存在，且页面上至少存在可编辑的"转盘动画"承载对象
  若维度 1 不满足 -> 直接 0 分

维度 2（完成度评分）：
    - 由若干 +x 加分细则和若干 -x 扣分细则组成
    - 加分细则：必须满足该细则下的全部要点才计入加分
    - 扣分细则：满足该细则下任意一点即计入扣分
    - 最终得分 = 各细则得分之和（允许为负数 -> 这里向下截断到 0，
      因为题目只说"累计"，未明确允许负总分；为安全起见仍打印原始累计）

说明
====
PPT 是一个静态描述格式，"动画行为"（旋转、按钮多次点击）无法在不实际播放
的情况下做严格自动判断。本脚本采用如下"灵活变通"策略：
    * 旋转/启停一致性：检查是否存在 PowerPoint 动画时间线（a:timing/p:timing
      或形状级动画属性），或是否使用了视频/媒体作为动画载体；若整页是单段
      视频，则按"视频动画"判断启停一致性（视频内部天然同步），并据此命中
      扣分项"整页被视频/动画覆盖，转盘和按钮无法单独编辑"。
    * "按钮再次点击/继续"：检测是否存在带触发器（trigger by click）的动画
      序列，且序列至少包含 start/pause/resume 三个语义状态（通过形状文字
      含"开始/暂停/继续"或同义词识别）。
"""

from __future__ import annotations

import json
import os
import re
import sys
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass, field
from typing import List, Tuple, Optional

from pptx import Presentation
from pptx.util import Emu

# ----------------------------- 工具 ----------------------------- #

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

EMU_PER_CM = 360000.0


def emu_to_cm(v: Optional[int]) -> float:
    if v is None:
        return 0.0
    return v / EMU_PER_CM


def shape_center_cm(shape) -> Tuple[float, float]:
    left = emu_to_cm(shape.left or 0)
    top = emu_to_cm(shape.top or 0)
    w = emu_to_cm(shape.width or 0)
    h = emu_to_cm(shape.height or 0)
    return left + w / 2.0, top + h / 2.0


def shape_size_cm(shape) -> Tuple[float, float]:
    return emu_to_cm(shape.width or 0), emu_to_cm(shape.height or 0)


def safe_text(shape) -> str:
    try:
        if shape.has_text_frame:
            return shape.text_frame.text or ""
    except Exception:
        pass
    return ""


# --------------------------- 报告对象 --------------------------- #


@dataclass
class Rule:
    code: str            # e.g. "+5-a", "-3-b"
    score: int           # 正/负
    title: str
    matched: bool = False
    notes: List[str] = field(default_factory=list)


@dataclass
class Report:
    dim1_pass: bool = False
    dim1_reasons: List[str] = field(default_factory=list)
    rules: List[Rule] = field(default_factory=list)

    def total(self) -> int:
        return sum(r.score for r in self.rules if r.matched)


# --------------------------- 主评估器 --------------------------- #


class Evaluator:
    TARGET_PAGE_INDEX = 9  # 第 10 页

    def __init__(self, path: str):
        self.path = path
        self.report = Report()
        self.prs: Optional[Presentation] = None
        self.slide = None
        self.slide_xml: Optional[ET.Element] = None
        self.slide_xml_raw: str = ""
        self.zf: Optional[zipfile.ZipFile] = None

    # -------- 维度 1 -------- #

    def check_dim1(self) -> bool:
        """维度 1：交付文件为 .pptx 格式，能够正常打开。
        严格按细则逐点：
            (1) 交付文件为 .pptx 格式
            (2) 能够正常打开
        """
        r = self.report

        # ---- (1) 后缀为 .pptx ---- #
        ext = os.path.splitext(self.path)[1].lower()
        if ext != ".pptx":
            r.dim1_reasons.append(f"(1)文件后缀 {ext} 非 .pptx")
            return False
        r.dim1_reasons.append(f"(1)文件后缀 {ext} 合法")

        # ---- (2) 能正常打开 ---- #
        # 代理判定：文件可被 python-pptx 完整加载，且包内 presentation.xml /
        # slideN.xml 可解析为合法 XML。
        try:
            self.prs = Presentation(self.path)
            self.zf = zipfile.ZipFile(self.path, "r")
            # 包合法性
            bad = self.zf.testzip()
            if bad is not None:
                r.dim1_reasons.append(f"(2)文件损坏，无法打开（坏文件：{bad}）")
                return False
        except Exception as e:
            r.dim1_reasons.append(f"(2)文件无法正常打开：{e}")
            return False
        r.dim1_reasons.append("(2)文件可正常打开")

        assert self.prs is not None
        if len(self.prs.slides) > self.TARGET_PAGE_INDEX:
            self.slide = self.prs.slides[self.TARGET_PAGE_INDEX]
        if self.slide is not None:
            self.slide_xml_raw = self.zf.read(
                f"ppt/slides/slide{self.TARGET_PAGE_INDEX + 1}.xml"
            ).decode("utf-8")
            self.slide_xml = ET.fromstring(self.slide_xml_raw)

        return True

    def _has_editable_spinner_object(self) -> bool:
        """第 10 页是否存在"可编辑的转盘动画对象"。
        判定：
          - 必须是 PPT 内可单独选中并编辑的 shape（OVAL/PICTURE/GROUP/FREEFORM 等），
            排除整页 MEDIA（视频）和整页 PICTURE 这类"烘焙背景"——它们不
            可被单独编辑为"转盘"，只是被整体播放/显示的烘焙内容。
          - 该对象近似为圆形（宽高差 ≤ 0.2cm，作为"轮盘"的几何代理）。
          - 该对象尺寸需要在合理的"转盘"量级：直径 ≥ 3cm（细则建议 6~7cm，
            这里给出较宽容差，但要排除明显是图标/指针的小图）。
          - 同时具有动画属性的代理：页面包含 p:timing 且时间轴中存在作用于
            该 shape 的旋转/动作动画；若仅有页面级 timing 但目标是其它对象
            （例如整页视频），则不视为"该对象本身的动画"。
        """
        if self.slide is None or self.slide_xml is None or self.prs is None:
            return False
        sw = emu_to_cm(self.prs.slide_width)
        sh = emu_to_cm(self.prs.slide_height)

        # 收集 timing 中所有动画目标 spid -> 用于"是否对该 shape 有动画"
        anim_target_ids = set()
        for spt in self.slide_xml.findall(".//p:spTgt", NS):
            sid = spt.get("spid")
            if sid:
                anim_target_ids.add(sid)

        def shape_id(sh):
            for tag in ("p:nvSpPr/p:cNvPr", "p:nvPicPr/p:cNvPr",
                        "p:nvGrpSpPr/p:cNvPr"):
                el = sh._element.find(".//" + tag, NS)
                if el is not None:
                    return el.get("id")
            return None

        for shape in self.slide.shapes:
            st = str(shape.shape_type)
            w, h = shape_size_cm(shape)
            # 排除整页媒体（视频）—— 这是"被烘焙的整页动画"，不是"可编辑转盘"
            if "MEDIA" in st:
                continue
            # 排除整页图片背景
            if "PICTURE" in st and w >= sw * 0.9 and h >= sh * 0.9:
                continue
            # 仅认可可作为转盘载体的 shape 类型
            if not any(k in st for k in ("OVAL", "PICTURE", "GROUP", "FREEFORM")):
                continue
            # 圆形几何代理（宽高差 ≤ 0.2cm）
            if w <= 0 or h <= 0 or abs(w - h) > 0.2:
                continue
            # 尺寸需达到"转盘"量级：直径 ≥ 3cm；过小则是图标/指针
            d = (w + h) / 2.0
            if d < 3.0:
                continue
            # 动画属性：要么页面 timing 有作用于该 shape 的动画，
            # 要么至少 shape 上挂了 hyperlink/action（可用作触发器入口）。
            # 若两者皆无：仍允许，但需要至少 timing 节点存在（页面已具备动画体系）。
            sid = shape_id(shape)
            has_anim_on_this = sid is not None and sid in anim_target_ids
            page_has_timing = self.slide_xml.find(".//p:timing", NS) is not None
            if not (has_anim_on_this or page_has_timing):
                continue
            return True

        return False

    # -------- 通用：抽取页面上的关键对象 -------- #

    def _classify_shapes(self):
        """把第 10 页 shape 分类，便于后续规则。"""
        info = {
            "all": [],
            "ovals": [],          # 圆/椭圆形状
            "pictures": [],
            "media": [],          # 视频/音频
            "groups": [],
            "text_shapes": [],
            "buttons": [],
            "fullpage_media": [],  # 覆盖整页的媒体
            "fullpage_picture": [],
            "cross_x": [],         # 中央乘号 ×
            "pointers": [],        # 指针（小三角/小图片）
        }
        sw_cm = emu_to_cm(self.prs.slide_width)
        sh_cm = emu_to_cm(self.prs.slide_height)
        for shape in self.slide.shapes:
            info["all"].append(shape)
            st = str(shape.shape_type)
            w_cm, h_cm = shape_size_cm(shape)
            txt = safe_text(shape).strip()
            if "MEDIA" in st:
                info["media"].append(shape)
                if w_cm >= sw_cm * 0.9 and h_cm >= sh_cm * 0.9:
                    info["fullpage_media"].append(shape)
            elif "PICTURE" in st:
                info["pictures"].append(shape)
                if w_cm >= sw_cm * 0.9 and h_cm >= sh_cm * 0.9:
                    info["fullpage_picture"].append(shape)
            elif "GROUP" in st:
                info["groups"].append(shape)
            elif "OVAL" in st:
                info["ovals"].append(shape)
            if txt:
                info["text_shapes"].append(shape)
                if re.search(r"[×xX✕✖]", txt) and len(txt) <= 3:
                    info["cross_x"].append(shape)
                if re.search(r"(开始|暂停|继续|停止|播放|start|pause|resume|stop|play)",
                             txt, re.I):
                    info["buttons"].append(shape)
            # 指针：尺寸较小且贴近"轮盘上方"的图片/三角；启发式留到具体规则里
        return info, sw_cm, sh_cm

    # -------- 维度 2 规则 -------- #

    def check_dim2(self):
        info, sw, sh = self._classify_shapes()

        # --- +5 转盘数量与对称分布 ---
        self._rule_two_spinners(info, sw, sh)
        # --- +1 中央乘号 ---
        self._rule_central_cross(info, sw, sh)
        # --- +3 共同控制按钮 ---
        self._rule_common_button(info, sw, sh)
        # --- +5 动画行为（同步启停 + 持续旋转 + 仅绕圆心） ---
        self._rule_anim_behavior(info)
        # --- +1 固定指针 ---
        self._rule_pointers(info, sw, sh)
        # --- +1 中央乘号格式 ---
        self._rule_cross_format(info)
        # --- +1 按钮状态文字 ---
        self._rule_button_states(info)

        # --- 扣分 ---
        self._rule_minus_fullpage_video(info)

    # ---------- 加分项实现 ---------- #

    def _rule_two_spinners(self, info, sw, sh):
        """+5 严格按细则逐点：
            (1) 左右各一个，共 2 个完整转盘（页面上"完整转盘"实际数量恰为 2）
            (2) 左侧转盘：圆形轮盘 + 多个彩色扇区 + 可辨认数字/口算
            (3) 右侧转盘：圆形轮盘 + 多个彩色扇区 + 可辨认数字/口算
            (4) 两转盘直径一致（建议约 6~7cm）
            (5) 每个转盘 宽高差 ≤ 0.2cm
            (6) 左右对称分布在粉色椭圆两侧（以粉色椭圆中心为对称轴）
            (7) 两个转盘中心点处于同一水平线上
        """
        rule = Rule("+5-spinners", 5, "页面左右各1个完整转盘(共2)，含扇区与数字，"
                                       + "直径一致(约6~7cm)，宽高差≤0.2，左右对称"
                                       + "(以粉色椭圆中心为对称轴)，中心同一水平线")

        # 整页视频存在时无法核验转盘内容，跳过不计分
        if info["fullpage_media"]:
            rule.notes.append("存在整页视频，无法核验转盘细则，跳过不计分")
            self.report.rules.append(rule)
            return

        # 候选：第10页上"可作为整体转盘"的对象。
        candidates = []
        for s in (info["ovals"] + info["pictures"] + info["groups"]):
            if s in info["fullpage_media"] or s in info["fullpage_picture"]:
                continue
            candidates.append(s)

        # ---- (5) 每个转盘自身 宽高差 ≤ 0.2cm（圆形轮盘的代理判定） ---- #
        def is_circular(s):
            w, h = shape_size_cm(s)
            return w > 0 and h > 0 and abs(w - h) <= 0.2

        # ---- "完整转盘"筛选：圆形 + 尺寸达到转盘量级（直径 ≥3cm，
        #      细则建议 6~7cm；上限 12cm 以排除装饰性大圆环/背景） ---- #
        def is_complete_spinner(s):
            if not is_circular(s):
                return False
            w, h = shape_size_cm(s)
            d = (w + h) / 2.0
            return 3.0 <= d <= 12.0

        circular = [s for s in candidates if is_circular(s)]
        complete = [s for s in candidates if is_complete_spinner(s)]
        rule.notes.append(
            f"候选(非整页){len(candidates)} 个，圆形(宽高差≤0.2){len(circular)} 个，"
            + f"完整转盘(圆形且直径3~12cm){len(complete)} 个"
        )

        # ---- (1) 完整转盘数量"恰好"为 2 ---- #
        # 与之前"从更多候选里挑最像的两个"不同：这里严格要求实际完整转盘数=2；
        # 多余的完整转盘（>2）也判为不达标，避免"多余转盘"误得分。
        ok_count = (len(complete) == 2)
        picked = complete if ok_count else []
        rule.notes.append(f"(1)完整转盘数量恰为2={ok_count}")

        # ---- (4) 两转盘直径一致（宽高差≤0.2cm；6~7cm 为建议值，不强制） ---- #
        ok_same_diameter = False
        if ok_count:
            w1, h1 = shape_size_cm(picked[0])
            w2, h2 = shape_size_cm(picked[1])
            d1 = (w1 + h1) / 2.0
            d2 = (w2 + h2) / 2.0
            ok_same_diameter = abs(d1 - d2) <= 0.2
            rule.notes.append(
                f"直径 D1={d1:.2f}cm D2={d2:.2f}cm；一致(|ΔD|≤0.2)={ok_same_diameter}（建议6~7cm）"
            )

        # ---- (6) 左右对称分布在"粉色椭圆两侧"（以椭圆中心为对称轴） ---- #
        # 严格按 rubric：对称轴是粉色椭圆中心的 x 坐标，而非幻灯片中线。
        # 找不到粉色椭圆时，本点判为不满足（不再退化到幻灯片中线）。
        # ---- (7) 中心点处于同一水平线 ---- #
        ok_symmetric = False
        ok_same_y = False
        left_shape = right_shape = None
        pink_xy = self._find_pink_ellipse_center(info, sw, sh)
        if ok_count:
            a, b = picked
            ca = shape_center_cm(a)
            cb = shape_center_cm(b)
            left_shape, right_shape = (a, b) if ca[0] < cb[0] else (b, a)
            cl = shape_center_cm(left_shape)
            cr = shape_center_cm(right_shape)
            if pink_xy is not None:
                axis_x = pink_xy[0]
                on_two_sides = cl[0] < axis_x < cr[0]
                sym_dist = abs((axis_x - cl[0]) - (cr[0] - axis_x)) <= 0.3
                ok_symmetric = on_two_sides and sym_dist
                rule.notes.append(
                    f"粉色椭圆中心 x={axis_x:.2f}cm；左中心={cl}, 右中心={cr}；"
                    f"分居两侧={on_two_sides}, 到椭圆中心距离对称(≤0.3cm)={sym_dist}"
                )
            else:
                rule.notes.append("未定位到粉色椭圆 -> (6)左右对称无法核验，判为不满足")
            # (7) 中心 y 在同一水平线，细则未给具体容差，按 0.3cm 工程容差判定
            ok_same_y = abs(cl[1] - cr[1]) <= 0.3
            rule.notes.append(
                f"中心 y 差={abs(cl[1]-cr[1]):.3f}cm 同一水平线(≤0.3cm)={ok_same_y}"
            )

        # ---- (2)(3) 含彩色扇区 + 可辨认数字/口算 ---- #
        # 检测策略：在 slide xml 中，对该 shape 子树内
        #   · 矢量承载：≥3 种不同的 srgbClr 填充色（扇区）+ shape 内文本含数字或口算符号
        #   · 图片承载：只有当图片的 alt 文本 / 图形名称 / 关系目标（media 文件名）
        #     明确显示"转盘/扇区/数字/口算"等语义时，才视为"可辨认"。仅仅是
        #     一张位图（无任何证据）不能被直接判为含扇区和含数字。
        def _iter_alt_and_names(shape):
            """收集 shape 上可用作图片语义证据的字符串："""
            texts = []
            # 图形名
            if shape.name:
                texts.append(shape.name)
                # descr / title on cNvPr
            for cnv_path in (".//p:nvSpPr/p:cNvPr",
                             ".//p:nvPicPr/p:cNvPr",
                             ".//p:nvGrpSpPr/p:cNvPr"):
                el = shape._element.find(cnv_path, NS)
                if el is not None:
                    for a in ("descr", "title", "name"):
                        v = el.get(a)
                        if v:
                            texts.append(v)
            # 图片承载对应 media 文件名（若能取到）
            try:
                if hasattr(shape, "image") and shape.image is not None:
                    fn = getattr(shape.image, "filename", "") or ""
                    if fn:
                        texts.append(fn)
            except Exception:
                pass
            return " ".join(texts)

        SECTOR_HINT = re.compile(
            r"(转盘|轮盘|扇区|扇形|spinner|wheel|sector|pie)", re.I
        )
        DIGIT_HINT = re.compile(
            r"(口算|算式|数字|number|digit|math|arith|equation|"
            r"[0-9０-９]|[\+\-\×\÷])", re.I
        )

        def has_sectors_and_digits(shape):
            xml = shape._element
            xml_str = ET.tostring(xml, encoding="unicode")
            # 矢量填充色种类（扇区代理）
            colors = set(re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', xml_str))
            many_colors = len(colors) >= 3
            # 矢量文字：数字/口算
            txt = safe_text(shape)
            has_digit_text = bool(re.search(r"[0-9０-９]", txt)) or \
                             bool(re.search(r"[\+\-\×\÷xX]", txt))
            # 图片承载
            is_image_carrier = re.search(r'r:embed="', xml_str) is not None or \
                               "PICTURE" in str(shape.shape_type)
            # 图片语义证据（alt / name / filename）
            alt_txt = _iter_alt_and_names(shape) if is_image_carrier else ""
            img_sector_evidence = bool(SECTOR_HINT.search(alt_txt))
            img_digit_evidence = bool(DIGIT_HINT.search(alt_txt))

            # 扇区：矢量多色 或 图片且有扇区语义证据
            sectors_ok = many_colors or (is_image_carrier and img_sector_evidence)
            # 数字：矢量文本含数字/口算 或 图片且有数字/口算语义证据
            digits_ok = has_digit_text or (is_image_carrier and img_digit_evidence)
            return sectors_ok, digits_ok, {
                "colors": len(colors),
                "is_image": is_image_carrier,
                "alt": alt_txt[:60],
                "text": txt[:30],
            }

        ok_left_content = ok_right_content = False
        if ok_count:
            s_ok_l, d_ok_l, info_l = has_sectors_and_digits(left_shape)
            s_ok_r, d_ok_r, info_r = has_sectors_and_digits(right_shape)
            ok_left_content = s_ok_l and d_ok_l
            ok_right_content = s_ok_r and d_ok_r
            rule.notes.append(
                f"左转盘：扇区={s_ok_l} 数字/口算={d_ok_l} 详情={info_l}"
            )
            rule.notes.append(
                f"右转盘：扇区={s_ok_r} 数字/口算={d_ok_r} 详情={info_r}"
            )

        # 圆形轮盘：(5) 已通过 is_circular 验证
        ok_circular_two = ok_count and is_circular(picked[0]) and is_circular(picked[1])

        rule.matched = (
            ok_count           # (1) 完整转盘数=2
            and ok_circular_two  # (2)(3) 圆形轮盘 + (5) 宽高差≤0.2cm
            and ok_left_content  # (2) 左：扇区+数字
            and ok_right_content # (3) 右：扇区+数字
            and ok_same_diameter # (4) 直径一致
            and ok_symmetric     # (6) 以粉色椭圆中心为对称轴左右对称
            and ok_same_y        # (7) 中心同一水平线
        )
        rule.notes.append(
            f"逐点：(1)完整转盘数=2 {ok_count} "
            f"(2/3)圆形={ok_circular_two} "
            f"左内容={ok_left_content} 右内容={ok_right_content} "
            f"(4)同径={ok_same_diameter} "
            f"(6)对称(粉椭圆轴)={ok_symmetric} (7)同水平={ok_same_y}"
        )
        self.report.rules.append(rule)

    def _rule_central_cross(self, info, sw, sh):
        """+1 严格按细则逐点：
            (1) 第 10 页存在"中央粉色椭圆"
            (2) 在该椭圆中心点放置"×"
            (3) × 中心相对椭圆中心的偏离不超过 0.3 厘米
        细则没有要求"× 是独立文本框"，因此只要页面上存在以
        "×/✕/✖/x/X" 为内容的对象（文本/小图片），都参与判定。
        """
        rule = Rule("+1-cross", 1, "中央粉色椭圆中心放置 ×，偏离≤0.3cm")

        # 整页视频存在时无法核验视频内部的粉色椭圆和 ×，跳过不计分
        if info["fullpage_media"]:
            rule.notes.append("存在整页视频，无法核验中央乘号细则，跳过不计分")
            self.report.rules.append(rule)
            return

        # ---- (1) 定位"中央粉色椭圆" ---- #
        # 方式 A：页面上有 OVAL/freeform 形状且填充为粉色（H ≈ 320~360°，S≥0.1，L≥0.7）
        # 方式 B：作为底图背景的一部分 -> 通过对背景图做像素级粉色连通域分析定位中心
        pink_ellipse_center = self._find_pink_ellipse_center(info, sw, sh)
        if pink_ellipse_center is None:
            rule.notes.append("未在第10页定位到中央粉色椭圆")
            self.report.rules.append(rule)
            return
        tx, ty = pink_ellipse_center
        rule.notes.append(f"中央粉色椭圆中心定位于 ({tx:.2f}cm, {ty:.2f}cm)")

        # ---- (2) 找 × 对象 ---- #
        cross_shapes = []
        for s in self.slide.shapes:
            if s in info["fullpage_media"] or s in info["fullpage_picture"]:
                continue
            t = safe_text(s).strip()
            if t and re.fullmatch(r"[×xX✕✖]+", t):
                cross_shapes.append(s)
        if not cross_shapes:
            rule.notes.append("未发现 × 对象")
            self.report.rules.append(rule)
            return

        # ---- (3) 偏离 ≤ 0.3cm ---- #
        for s in cross_shapes:
            cx, cy = shape_center_cm(s)
            dist = ((cx - tx) ** 2 + (cy - ty) ** 2) ** 0.5
            rule.notes.append(
                f"× '{safe_text(s).strip()}' 中心({cx:.2f},{cy:.2f})，"
                f"距椭圆中心 {dist:.3f}cm"
            )
            if dist <= 0.3:
                rule.matched = True
                break
        self.report.rules.append(rule)

    def _find_pink_ellipse_center(self, info, sw, sh):
        """返回中央粉色椭圆的中心 (x_cm, y_cm)；找不到返回 None。
        优先：页面上"形状级"粉色椭圆 -> 取该 shape 中心。
        其次：解析背景图（整页图片）中粉色连通域，取最居中的连通域几何中心，
              再换算回幻灯片厘米坐标。
        """
        # ---- A) 形状级粉色椭圆 ---- #
        for s in self.slide.shapes:
            if s in info["fullpage_media"] or s in info["fullpage_picture"]:
                continue
            st = str(s.shape_type)
            if "OVAL" not in st and "FREEFORM" not in st:
                continue
            xml_str = ET.tostring(s._element, encoding="unicode")
            for hex6 in re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', xml_str):
                if self._is_pink(hex6):
                    cx, cy = shape_center_cm(s)
                    return (cx, cy)

        # ---- B) 背景图像素分析 ---- #
        bg = None
        for s in self.slide.shapes:
            if s in info["fullpage_picture"]:
                bg = s
                break
        if bg is None:
            return None
        try:
            from PIL import Image
            import io
            blob = bg.image.blob
            img = Image.open(io.BytesIO(blob)).convert("RGB")
        except Exception:
            return None

        W, H = img.size
        px = img.load()
        # 在中央带状区域内搜索粉色像素，避免边角粉色装饰干扰
        x0, x1 = int(W * 0.35), int(W * 0.65)
        y0, y1 = int(H * 0.30), int(H * 0.70)
        xs, ys, n = 0, 0, 0
        for yy in range(y0, y1, 2):
            for xx in range(x0, x1, 2):
                r, g, b = px[xx, yy]
                if self._is_pink_rgb(r, g, b):
                    xs += xx
                    ys += yy
                    n += 1
        if n < 50:
            return None
        cx_px = xs / n
        cy_px = ys / n
        # 像素 -> 幻灯片厘米
        bg_left = emu_to_cm(bg.left or 0)
        bg_top = emu_to_cm(bg.top or 0)
        bg_w = emu_to_cm(bg.width or 0)
        bg_h = emu_to_cm(bg.height or 0)
        cx_cm = bg_left + bg_w * (cx_px / W)
        cy_cm = bg_top + bg_h * (cy_px / H)
        return (cx_cm, cy_cm)

    @staticmethod
    def _is_pink(hex6: str) -> bool:
        try:
            r = int(hex6[0:2], 16)
            g = int(hex6[2:4], 16)
            b = int(hex6[4:6], 16)
        except Exception:
            return False
        return Evaluator._is_pink_rgb(r, g, b)

    @staticmethod
    def _is_pink_rgb(r: int, g: int, b: int) -> bool:
        # 粉色：红明显高，蓝中高，绿略低；亮度较高
        if r < 200:
            return False
        if not (b >= 150 and b <= 240):
            return False
        if not (g >= 130 and g <= 220):
            return False
        if not (r > g and r >= b - 10):
            return False
        # 排除纯白/灰
        if abs(r - g) < 15 and abs(g - b) < 15:
            return False
        return True


    def _rule_common_button(self, info, sw, sh):
        """+3 严格按细则逐点：
            (1) 页面中"仅"设置一个主要按钮
            (2) 该按钮"同时控制"左、右两个转盘
            (3) 位于"两个转盘之间" 或 "中央粉色椭圆下方"
            (4) 宽约 3~5 厘米
            (5) 高约 1~1.5 厘米
        """
        rule = Rule("+3-common-button", 3,
                    "页面仅1个主要按钮，同控两转盘，位于两转盘之间或粉椭圆下方，"
                    "宽3~5cm 高1~1.5cm")

        # ---- 收集"按钮候选" ---- #
        # 细则没有限定按钮的实现形式（可以是 shape/图片/文本框/group/媒体控件），
        # 因此把页面上"主要可交互对象"都列为候选：
        # 排除：整页媒体/整页背景图、纯装饰的转盘/指针/× 等。
        # 主要按钮的可识别特征（任一即可作为候选）：
        #   a) 形状文字命中"开始/暂停/继续/停止/播放/start/pause/resume/stop/play"
        #   b) 在 slide xml 中绑定了 hyperlink/action(ppaction://) -> 可点击
        #   c) shape 名称含 "button/btn/控制/按钮"

        def is_button_like(sh):
            # a) 文字
            t = safe_text(sh)
            if t and re.search(r"(开始|暂停|继续|停止|播放|start|pause|resume|stop|play)",
                               t, re.I):
                return True
            # c) 名称
            name = (sh.name or "")
            if re.search(r"button|btn|控制|按钮", name, re.I):
                return True
            # b) 在该 shape 的子树 xml 中查找点击动作
            try:
                xml = ET.tostring(sh._element, encoding="unicode")
                if "ppaction://" in xml or "hlinkClick" in xml:
                    return True
            except Exception:
                pass
            return False

        buttons = []
        for sh in self.slide.shapes:
            if sh in info["fullpage_media"] or sh in info["fullpage_picture"]:
                continue
            if is_button_like(sh):
                buttons.append(sh)
        rule.notes.append(f"页面按钮候选数：{len(buttons)}")

        # 整页视频存在且没有找到独立按钮，跳过不计分
        if info["fullpage_media"] and not buttons:
            rule.notes.append("存在整页视频且页面无独立按钮，跳过不计分")
            self.report.rules.append(rule)
            return

        # ---- (1) 仅 1 个主要按钮 ---- #
        ok_only_one = len(buttons) == 1
        if not ok_only_one:
            rule.notes.append("不满足'仅1个主要按钮'")
        if not buttons:
            self.report.rules.append(rule)
            return

        btn = buttons[0]
        w, h = shape_size_cm(btn)
        cx, cy = shape_center_cm(btn)
        rule.notes.append(
            f"按钮：name={btn.name!r} 文字={safe_text(btn).strip()!r} "
            f"尺寸 {w:.2f}x{h:.2f}cm 中心({cx:.2f},{cy:.2f})"
        )

        # ---- (4)(5) 尺寸 ---- #
        ok_w = 3.0 <= w <= 5.0
        ok_h = 1.0 <= h <= 1.5
        rule.notes.append(f"(4)宽3~5cm={ok_w}  (5)高1~1.5cm={ok_h}")

        # ---- 先定位两个转盘（与 +5-spinners 保持一致：圆形 + 直径3~12cm，
        #      按接近 6.5cm 排序取前两个），供 (2) 目标一致性 与 (3) 位置判定共用 ---- #
        sp_cands = []
        for s in (info["ovals"] + info["pictures"] + info["groups"]):
            if s in info["fullpage_media"] or s in info["fullpage_picture"]:
                continue
            ww, hh = shape_size_cm(s)
            if ww <= 0 or hh <= 0 or abs(ww - hh) > 0.2:
                continue
            dd = (ww + hh) / 2.0
            if not (3.0 <= dd <= 12.0):
                continue
            sp_cands.append(s)
        sp_cands.sort(key=lambda s: abs((shape_size_cm(s)[0] +
                                         shape_size_cm(s)[1]) / 2 - 6.5))
        two_spinners = sp_cands[:2]

        def _shape_id(sh):
            for tag in ("p:nvSpPr/p:cNvPr", "p:nvPicPr/p:cNvPr",
                        "p:nvGrpSpPr/p:cNvPr"):
                el = sh._element.find(".//" + tag, NS)
                if el is not None:
                    return el.get("id")
            return None

        spinner_ids = {sid for sid in (_shape_id(s) for s in two_spinners) if sid}
        rule.notes.append(f"两个转盘 spid={spinner_ids}")

        # ---- (2) 同时控制左、右两个转盘 ---- #
        # 严格化：按钮 click 触发的目标 spid 集合必须"包含"两个转盘 spid，
        # 而不是仅"存在 ≥2 个不同的 target"。
        if len(spinner_ids) == 2:
            ok_control_both = self._button_controls_target_set(btn, spinner_ids)
        else:
            ok_control_both = False
            rule.notes.append("未定位到两个转盘 spid，(2)同控两转盘判为不满足")
        rule.notes.append(f"(2)按钮触发目标包含两个转盘 spid={ok_control_both}")

        # ---- (3) 位置：两个转盘之间 或 中央粉色椭圆下方 ---- #
        # 严格化：
        #   "两转盘之间" 需同时约束 x/y：
        #     · x 在两转盘 x 区间内
        #     · y 与两转盘中心 y 均值的偏差 ≤ 转盘半径 + 1.0cm 的工程容差
        #       （避免按钮位于页面很上方/很下方也被判为"之间"）
        #   "中央粉色椭圆下方" 需要 y 在椭圆中心之下 且 x 与椭圆中心同列
        pink_xy = self._find_pink_ellipse_center(info, sw, sh)
        ok_pos_between = False
        ok_pos_below_pink = False
        if len(two_spinners) == 2:
            cL = shape_center_cm(two_spinners[0])
            cR = shape_center_cm(two_spinners[1])
            xl, xr = min(cL[0], cR[0]), max(cL[0], cR[0])
            avg_y = (cL[1] + cR[1]) / 2.0
            # 转盘半径（用较大者，容差更宽容）
            radii = [(shape_size_cm(s)[0] + shape_size_cm(s)[1]) / 4.0
                     for s in two_spinners]
            r_max = max(radii) if radii else 0.0
            y_tol = r_max + 1.0
            x_between = xl <= cx <= xr
            y_between = abs(cy - avg_y) <= y_tol
            ok_pos_between = x_between and y_between
            rule.notes.append(
                f"两转盘中心 L={cL} R={cR}，平均y={avg_y:.2f}, 半径≈{r_max:.2f}cm, "
                + f"y容差={y_tol:.2f}cm；按钮中心=({cx:.2f},{cy:.2f}) "
                + f"x在[{xl:.2f},{xr:.2f}]内={x_between}, |cy-avg_y|≤容差={y_between}"
            )
        else:
            rule.notes.append("未找到两个转盘 -> 无法判定'两转盘之间'")
        if pink_xy is not None:
            px, py = pink_xy
            # "中央粉色椭圆下方"：按钮中心 y > 椭圆中心 y；并且 x 大致与椭圆同列
            ok_pos_below_pink = (cy > py) and (abs(cx - px) <= 3.0)
            rule.notes.append(
                f"粉椭圆中心=({px:.2f},{py:.2f})；按钮中心 y={cy:.2f} 在其下方且 x 同列={ok_pos_below_pink}"
            )
        else:
            rule.notes.append("未定位到粉色椭圆 -> 无法判定'椭圆下方'")
        ok_pos = ok_pos_between or ok_pos_below_pink
        rule.notes.append(f"(3)位置 之间或下方={ok_pos}")

        rule.matched = ok_only_one and ok_control_both and ok_pos and ok_w and ok_h
        rule.notes.append(
            f"逐点：(1)仅1个={ok_only_one} (2)同控两转盘={ok_control_both} "
            f"(3)位置={ok_pos} (4)宽={ok_w} (5)高={ok_h}"
        )
        self.report.rules.append(rule)

    def _button_controls_two_targets(self, btn) -> bool:
        """在 slide xml 的 timing 节里，按钮的 spid 作为 click trigger 是否
        驱动了 ≥2 个不同 spid 的目标。
        """
        targets = self._button_click_target_spids(btn)
        return len(targets) >= 2

    def _button_controls_target_set(self, btn, required_ids: "set[str]") -> bool:
        """按钮 click 触发的目标 spid 集合是否"包含"给定 spid 集合。
        与 _button_controls_two_targets 的差别：这里要求目标必须恰好覆盖
        指定的两个转盘 spid，而不仅仅是"至少 2 个不同 target"。
        """
        if not required_ids:
            return False
        targets = self._button_click_target_spids(btn)
        return required_ids.issubset(targets)

    def _button_click_target_spids(self, btn) -> "set[str]":
        """收集：以 btn 的 spid 为 onClick trigger 的所有 par 下动画 targetElement spid。"""
        # 按钮 spid
        try:
            spid = btn._element.find(".//p:nvSpPr/p:cNvPr", NS)
            if spid is None:
                spid = btn._element.find(".//p:nvPicPr/p:cNvPr", NS)
            if spid is None:
                spid = btn._element.find(".//p:nvGrpSpPr/p:cNvPr", NS)
            btn_id = spid.get("id") if spid is not None else None
        except Exception:
            btn_id = None
        targets: "set[str]" = set()
        if btn_id is None or self.slide_xml is None:
            return targets

        timing = self.slide_xml.find(".//p:timing", NS)
        if timing is None:
            return targets

        for par in timing.iter("{%s}par" % NS["p"]):
            ctn = par.find("p:cTn", NS)
            if ctn is None:
                continue
            if ctn.get("nodeType") not in ("clickEffect", "withEffect",
                                            "afterEffect", "interactiveSeq"):
                continue
            trigger_ok = False
            for cond in par.findall(".//p:cond", NS):
                if cond.get("evt") == "onClick":
                    sptgt = cond.find(".//p:spTgt", NS)
                    if sptgt is not None and sptgt.get("spid") == btn_id:
                        trigger_ok = True
                        break
            if not trigger_ok:
                seq = par.find(".//p:tgtEl/p:spTgt", NS)
                if seq is not None and seq.get("spid") == btn_id:
                    trigger_ok = True
            if not trigger_ok:
                continue
            for spt in par.findall(".//p:spTgt", NS):
                tid = spt.get("spid")
                if tid and tid != btn_id:
                    targets.add(tid)
        return targets

    def _rule_anim_behavior(self, info):
        """+5 严格按细则逐点：
            (1) 进入第 10 页后两个转盘保持静止，等待用户点击按钮
            (2) 按钮首次点击 -> 左、右两个转盘同时开始旋转，启动时间差 ≤ 0.2s
            (3) 两个转盘能够连续旋转，不会自动停止
            (4) 按钮再次点击 -> 两个转盘同时暂停，暂停时间差 ≤ 0.2s
            (5) 再次点击同一按钮 -> 两个转盘能够从暂停状态继续旋转
            (6) 两个转盘的"开始/暂停/继续"时刻一致，不出现一动一静
            (7) 旋转时仅轮盘绕各自圆心转动；转盘整体不发生位移、跳动、
                缩放或明显偏心
        所有点必须同时满足才计 +5。

        重要原则：本项属于交互播放行为，无法完全静态确认。
          · 若按钮绑定 VBA 宏 / 自定义 add-in / 未知 ppaction:// 目标 ->
            无法静态确认其行为，(2)(3)(4)(5)(6) 均按"不确定"处理，
            规则不命中（保守判定），并在 notes 中说明原因。
          · 只有当 p:timing 中给出结构化的 onClick + animRot / cmd
            证据，且证据同时覆盖两个转盘 spid，才能确认对应细则。
        """
        rule = Rule("+5-anim", 5,
                    "进入静止/同步启动(≤0.2s)/持续旋转/同步暂停(≤0.2s)/可继续/"
                    + "时刻一致/仅绕圆心（无法静态确认时保守不命中）")

        if self.slide_xml is None:
            rule.notes.append("无法读取第10页 XML")
            self.report.rules.append(rule)
            return
        timing = self.slide_xml.find(".//p:timing", NS)

        # ---- 找两个转盘（与 +5-spinners 一致：圆形 + 直径 3~12cm，
        #      按接近 6.5cm 取前两个） ---- #
        cands = []
        for s in (info["ovals"] + info["pictures"] + info["groups"]):
            if s in info["fullpage_media"] or s in info["fullpage_picture"]:
                continue
            w, h = shape_size_cm(s)
            if w <= 0 or h <= 0 or abs(w - h) > 0.2:
                continue
            d = (w + h) / 2.0
            if not (3.0 <= d <= 12.0):
                continue
            cands.append(s)
        cands.sort(key=lambda s: abs((shape_size_cm(s)[0] +
                                      shape_size_cm(s)[1]) / 2 - 6.5))
        spinners = cands[:2]

        if len(spinners) < 2:
            rule.notes.append(f"页面找不到 2 个完整转盘 -> 无法核对动画(找到 {len(spinners)})")
            self.report.rules.append(rule)
            return

        if timing is None:
            rule.notes.append("不存在 p:timing 动画时间轴 -> 无法静态确认交互行为，判为不命中")
            self.report.rules.append(rule)
            return

        def shape_id(sh):
            for tag in ("p:nvSpPr/p:cNvPr", "p:nvPicPr/p:cNvPr",
                        "p:nvGrpSpPr/p:cNvPr"):
                el = sh._element.find(".//" + tag, NS)
                if el is not None:
                    return el.get("id")
            return None

        sp_ids = [shape_id(s) for s in spinners]
        sp_id_set = {i for i in sp_ids if i}
        rule.notes.append(f"两转盘 spid={sp_ids}")
        if len(sp_id_set) < 2:
            rule.notes.append("两个转盘的 spid 无法唯一定位 -> 判为不命中")
            self.report.rules.append(rule)
            return

        # ---- 找"主要按钮"及其交互形式 ---- #
        btn = None
        for sh in self.slide.shapes:
            if sh in info["fullpage_media"] or sh in info["fullpage_picture"]:
                continue
            t = safe_text(sh)
            name = sh.name or ""
            try:
                xml = ET.tostring(sh._element, encoding="unicode")
            except Exception:
                xml = ""
            if (t and re.search(r"(开始|暂停|继续|停止|播放|start|pause|resume|stop|play)", t, re.I)) \
               or re.search(r"button|btn|控制|按钮", name, re.I) \
               or "ppaction://" in xml or "hlinkClick" in xml:
                btn = sh
                break
        btn_id = shape_id(btn) if btn is not None else None
        rule.notes.append(f"主要按钮 spid={btn_id}")
        if btn is None or btn_id is None:
            rule.notes.append("找不到主要按钮 -> 判为不命中")
            self.report.rules.append(rule)
            return

        # ---- 若按钮绑定 VBA 宏 / 自定义动作，则无法静态确认交互行为 ---- #
        # 判据：按钮子树包含 <a:hlinkClick action="ppaction://macro?..."> 或
        # 演示文稿中存在 ppt/vbaProject.bin 且按钮的 hlinkClick action 值不是
        # 我们能静态识别的形式（如 program=... / hlinkshowjump=...）
        try:
            btn_xml = ET.tostring(btn._element, encoding="unicode")
        except Exception:
            btn_xml = ""
        macro_action = re.search(r'action="ppaction://macro[^"]*"', btn_xml)
        has_vba_project = False
        try:
            if self.zf is not None and "ppt/vbaProject.bin" in self.zf.namelist():
                has_vba_project = True
        except Exception:
            pass
        if macro_action or (has_vba_project and re.search(r'action="ppaction://program', btn_xml)):
            rule.notes.append(
                "按钮绑定 VBA 宏/自定义程序动作 -> 无法静态确认交互播放行为，"
                + "(2)(3)(4)(5)(6) 均按不确定处理，规则保守判为不命中"
            )
            self.report.rules.append(rule)
            return

        # ---- 收集 timing 中由该按钮 onClick 触发的 par ---- #
        click_pars = []
        for par in timing.iter("{%s}par" % NS["p"]):
            ctn = par.find("p:cTn", NS)
            if ctn is None:
                continue
            trig_btn = False
            for cond in par.findall(".//p:cond", NS):
                if cond.get("evt") == "onClick":
                    sptgt = cond.find(".//p:spTgt", NS)
                    if sptgt is not None and sptgt.get("spid") == btn_id:
                        trig_btn = True
                        break
            if not trig_btn:
                continue
            click_pars.append(par)
        rule.notes.append(f"由按钮 onClick 触发的 par 数={len(click_pars)}")

        # ---- 辅助：给定 par，统计其下 animRot 对两个转盘的作用 ---- #
        def rot_info_for_par(par):
            """返回 { spid: [(delay_ms, dur_ms, by_val, repeat)] }，只统计
            作用于转盘 spid 的 animRot 记录。"""
            m = {sid: [] for sid in sp_id_set}
            for ar in par.findall(".//p:animRot", NS):
                # 目标 spid
                tgt_ids = [spt.get("spid") for spt in ar.findall(".//p:spTgt", NS)]
                if not any(t in sp_id_set for t in tgt_ids):
                    continue
                ctn = ar.find(".//p:cTn", NS)
                delay_ms = 0
                dur_ms = 0
                repeat = "1"
                if ctn is not None:
                    d = ctn.get("delay", "0")
                    if d.isdigit():
                        delay_ms = int(d)
                    du = ctn.get("dur", "0")
                    if du.isdigit():
                        dur_ms = int(du)
                    repeat = ctn.get("repeatCount", "1")
                by_val = ar.get("by", "")  # animRot 的 by 单位为 1/60000 度
                for t in tgt_ids:
                    if t in sp_id_set:
                        m[t].append((delay_ms, dur_ms, by_val, repeat))
            return m

        # ---- 辅助：给定 par，统计其下"暂停/继续"命令对两转盘的作用 ---- #
        # 兼容 PowerPoint 常见编码：<p:cmd type="evt" cmd="togglePause"/>
        # 以及 <p:cmd cmd="resume"/> <p:cmd cmd="pause"/>
        def cmd_info_for_par(par):
            pause_targets = set()
            resume_targets = set()
            toggle_targets = set()
            for cmd in par.findall(".//p:cmd", NS):
                cmd_val = (cmd.get("cmd", "") or (cmd.text or "")).strip()
                lc = cmd_val.lower()
                if not cmd_val:
                    continue
                spids = [spt.get("spid") for spt in cmd.findall(".//p:spTgt", NS)]
                if "togglepause" in lc:
                    for sid in spids:
                        if sid:
                            toggle_targets.add(sid)
                elif "pause" in lc or "stop" in lc:
                    for sid in spids:
                        if sid:
                            pause_targets.add(sid)
                elif "resume" in lc or "play" in lc:
                    for sid in spids:
                        if sid:
                            resume_targets.add(sid)
            return pause_targets, resume_targets, toggle_targets

        # ---- (2) 首次点击 -> 同步启动，delay 差 ≤ 200ms ---- #
        ok_sync_start = False
        first_par = None
        first_rot_info = None
        for par in click_pars:
            rot_map = rot_info_for_par(par)
            if not all(rot_map[sid] for sid in sp_id_set):
                continue
            # 每个转盘取首条 animRot 的 delay
            first_delays = [rot_map[sid][0][0] for sid in sp_id_set]
            diff_ms = max(first_delays) - min(first_delays)
            if diff_ms <= 200:
                ok_sync_start = True
                first_par = par
                first_rot_info = rot_map
                rule.notes.append(
                    f"(2)首次点击 par 同步启动两转盘 animRot, delays={first_delays}ms, 差={diff_ms}ms"
                )
                break
        if not ok_sync_start:
            rule.notes.append("(2)未找到 onClick+animRot 同时覆盖两转盘、delay 差≤200ms 的 par")

        # ---- (3) 连续旋转：repeatCount=indefinite 且 by 表示至少一整圈 ---- #
        # animRot by 单位为 1/60000 度；一整圈 = 360 * 60000 = 21600000
        ok_continuous = False
        if first_par is not None and first_rot_info is not None:
            per_spinner_ok = []
            details = []
            for sid in sp_id_set:
                recs = first_rot_info[sid]
                any_ok = False
                for (delay_ms, dur_ms, by_val, repeat) in recs:
                    rep_ok = (repeat in ("indefinite", "-1", "INF")
                              or (repeat.isdigit() and int(repeat) >= 100000))
                    # by 若为空(默认或未指定)也视为可接受，只要 repeat 无限；
                    # 若给出 by 数值，要求 |by| >= 一整圈（避免只旋转一点点）
                    by_ok = True
                    if by_val:
                        try:
                            by_ok = abs(int(by_val)) >= 21600000
                        except Exception:
                            # 非整数：可能是浮点度数；粗略判定 >= 360
                            try:
                                by_ok = abs(float(by_val)) >= 360.0
                            except Exception:
                                by_ok = True
                    if rep_ok and by_ok:
                        any_ok = True
                        details.append(
                            f"spid={sid} repeat={repeat} by={by_val or '(默认)'}"
                        )
                        break
                per_spinner_ok.append(any_ok)
            ok_continuous = all(per_spinner_ok) and len(per_spinner_ok) == 2
            rule.notes.append(
                f"(3)连续旋转={ok_continuous}  " + "; ".join(details)
            )
        else:
            rule.notes.append("(3)无法判定(缺首点击 par)")

        # ---- (4) 再次点击 -> 同步暂停，delay 差 ≤ 200ms ---- #
        # 有效编码：另一个 onClick par 下含 pause 或 togglePause 命令，
        # 且命令目标覆盖两个转盘 spid。
        ok_sync_pause = False
        pause_par = None
        for par in click_pars:
            if par is first_par:
                continue
            pause_t, _resume_t, toggle_t = cmd_info_for_par(par)
            pause_like = pause_t | toggle_t
            if not sp_id_set.issubset(pause_like):
                continue
            # 收集 pause/toggle 命令的 delay
            delays = []
            for cmd in par.findall(".//p:cmd", NS):
                cmd_val = (cmd.get("cmd", "") or (cmd.text or "")).strip().lower()
                if not ("pause" in cmd_val or "togglepause" in cmd_val or "stop" in cmd_val):
                    continue
                spids = [spt.get("spid") for spt in cmd.findall(".//p:spTgt", NS)]
                if not any(sid in sp_id_set for sid in spids):
                    continue
                ctn = cmd.find(".//p:cTn", NS)
                d = ctn.get("delay", "0") if ctn is not None else "0"
                try:
                    delays.append(int(d) if d.isdigit() else 0)
                except Exception:
                    delays.append(0)
            diff_ms = (max(delays) - min(delays)) if delays else 0
            if diff_ms <= 200:
                ok_sync_pause = True
                pause_par = par
                rule.notes.append(f"(4)第二次点击 par 同步暂停两转盘, 差={diff_ms}ms")
                break
        if not ok_sync_pause:
            rule.notes.append("(4)未找到 onClick + (pause/togglePause) 同时覆盖两转盘、差≤200ms 的 par")

        # ---- (5) 再次点击 -> 继续旋转 ---- #
        # 有效编码：
        #   A. 存在另一个 onClick par，其 resume/play 命令目标覆盖两转盘；
        #   B. 存在同一 togglePause 序列（暂停 par 的命令类型为 togglePause，
        #      再次点击天然切换回旋转），此时视为"编码上支持继续"。
        ok_resume = False
        for par in click_pars:
            if par is first_par or par is pause_par:
                continue
            _p_t, resume_t, toggle_t = cmd_info_for_par(par)
            if sp_id_set.issubset(resume_t) or sp_id_set.issubset(toggle_t):
                ok_resume = True
                rule.notes.append("(5)发现独立的 resume/togglePause par 作用于两转盘")
                break
        if not ok_resume and pause_par is not None:
            _p_t, _r_t, toggle_t = cmd_info_for_par(pause_par)
            if sp_id_set.issubset(toggle_t):
                ok_resume = True
                rule.notes.append("(5)暂停 par 使用 togglePause 命令，天然支持暂停/继续切换")
        if not ok_resume:
            rule.notes.append("(5)未发现能静态确认的'从暂停继续'动画命令")

        # ---- (1) 进入静止：不存在"进入即播"的 animRot 作用于转盘 ---- #
        ok_idle_on_enter = True
        for par in timing.iter("{%s}par" % NS["p"]):
            ctn = par.find("p:cTn", NS)
            if ctn is None:
                continue
            delay = ctn.get("delay", "")
            node_type = ctn.get("nodeType", "")
            has_click = any(
                cond.get("evt") == "onClick"
                for cond in par.findall(".//p:cond", NS)
            )
            has_rot_on_spinner = any(
                spt.get("spid") in sp_id_set
                for ar in par.findall(".//p:animRot", NS)
                for spt in ar.findall(".//p:spTgt", NS)
            )
            if has_rot_on_spinner and not has_click and \
               (delay in ("", "0") or node_type in ("mainSeq", "tmRoot")):
                ok_idle_on_enter = False
                break
        rule.notes.append(f"(1)进入页面后转盘静止={ok_idle_on_enter}")

        # ---- (6) 时刻一致：(2)(4) 都满足同步差 ≤200ms 视为一致 ---- #
        ok_concurrent = ok_sync_start and ok_sync_pause
        rule.notes.append(f"(6)开始/暂停/继续时刻一致={ok_concurrent}")

        # ---- (7) 仅绕圆心、不位移/缩放/偏心 ---- #
        # 正向证据：两个转盘的 animRot 存在（(2) 已确认）；bbox 宽高差 ≤0.2cm
        #           （PowerPoint animRot 默认绕 bbox 中心旋转，因此圆形 bbox
        #            即可保证绕圆心）；shape 上无 animMotion / animScale。
        # 负向证据：不存在 animMotion / animScale 作用于转盘 spid；
        #           不存在 <a:xfrm> 中 offset 与 ext 相对错位（PPT 中 shape
        #            旋转对齐 bbox，无法通过静态 XML 显式表示"跳动/偏心"）。
        no_motion = True
        no_scale = True
        for ar in self.slide_xml.findall(".//p:animMotion", NS):
            for spt in ar.findall(".//p:spTgt", NS):
                if spt.get("spid") in sp_id_set:
                    no_motion = False
                    break
        for ar in self.slide_xml.findall(".//p:animScale", NS):
            for spt in ar.findall(".//p:spTgt", NS):
                if spt.get("spid") in sp_id_set:
                    no_scale = False
                    break
        no_eccentric = True
        for sh in spinners:
            w, h = shape_size_cm(sh)
            if abs(w - h) > 0.2:
                no_eccentric = False
                break
        # 正向：两转盘都有 animRot 作为旋转载体
        has_rot_both = ok_sync_start  # (2) 已经确认覆盖两个 spid
        ok_rotate_only = no_motion and no_scale and no_eccentric and has_rot_both
        rule.notes.append(
            f"(7)有animRot覆盖两转盘={has_rot_both} 无位移={no_motion} "
            + f"无缩放={no_scale} 无明显偏心={no_eccentric} -> 仅绕圆心={ok_rotate_only}"
        )

        rule.matched = (
            ok_idle_on_enter and ok_sync_start and ok_continuous and
            ok_sync_pause and ok_resume and ok_concurrent and ok_rotate_only
        )
        rule.notes.append(
            f"逐点：(1)={ok_idle_on_enter} (2)={ok_sync_start} (3)={ok_continuous} "
            + f"(4)={ok_sync_pause} (5)={ok_resume} (6)={ok_concurrent} (7)={ok_rotate_only}"
        )
        self.report.rules.append(rule)

    def _rule_pointers(self, info, sw, sh):
        """+1 两转盘上方各一个"指向轮盘边缘"的指针，大小和位置一致。

        严格化判定（回应第一轮反馈"缺少指向轮盘边缘的实现"）：
          - 指针形状线索：三角形/箭头/Freeform/自选图形中含
            "Triangle/Arrow/pointer/箭头/指针/三角"；单纯装饰图片不算。
          - 朝向：形状 bbox 宽/高比接近细长（长边:短边 ≥ 1.4），或旋转角度
            指示"向下朝向轮盘"；对三角形/箭头预设几何来推出尖端。
          - 尖端位置：估算尖端在页面上的坐标（cm），要求：
              · 尖端 y 位于转盘顶边附近（|y_tip - (sy - r)| ≤ 半径 * 0.35）
              · 尖端 x 与转盘中心 x 的偏差 ≤ 半径 * 0.35
            两条同时满足视为"指向轮盘边缘"。
          - 大小/位置一致仍要求两个指针宽高差 ≤0.2cm、相对各自转盘中心的
            偏移量差 ≤0.3cm。
        """
        rule = Rule("+1-pointers", 1,
                    "两转盘上方指针(形状/朝向/尖端在轮盘边缘)，大小位置一致")

        if self.slide is None:
            self.report.rules.append(rule)
            return

        # ---- 收集候选（排除整页媒体/图片） ---- #
        pointer_cands = []
        for s in self.slide.shapes:
            if s in info["fullpage_media"] or s in info["fullpage_picture"]:
                continue
            w, h = shape_size_cm(s)
            if w > 0 and h > 0:
                pointer_cands.append(s)

        if info["fullpage_media"] and not pointer_cands:
            rule.notes.append("存在整页视频且页面无独立对象，跳过不计分")
            self.report.rules.append(rule)
            return

        # ---- 找两个转盘（与其他规则一致：圆形 + 直径3~12cm） ---- #
        sp_cands = []
        for s in (info["ovals"] + info["pictures"] + info["groups"]):
            if s in info["fullpage_media"] or s in info["fullpage_picture"]:
                continue
            ww, hh = shape_size_cm(s)
            if ww <= 0 or hh <= 0 or abs(ww - hh) > 0.2:
                continue
            dd = (ww + hh) / 2.0
            if not (3.0 <= dd <= 12.0):
                continue
            sp_cands.append(s)
        sp_cands.sort(key=lambda s: abs((shape_size_cm(s)[0] +
                                         shape_size_cm(s)[1]) / 2 - 6.5))
        spinners = sp_cands[:2]

        if len(spinners) < 2:
            rule.notes.append("找不到两个转盘，无法判定指针位置")
            self.report.rules.append(rule)
            return

        spinners.sort(key=lambda s: shape_center_cm(s)[0])
        left_sp, right_sp = spinners
        cl_sp = shape_center_cm(left_sp)
        cr_sp = shape_center_cm(right_sp)
        l_rad = shape_size_cm(left_sp)[0] / 2.0
        r_rad = shape_size_cm(right_sp)[0] / 2.0

        # ---- 指针形状/朝向识别 ---- #
        POINTER_SHAPE_HINT = re.compile(
            r"(triangle|arrow|pointer|三角|箭头|指针)", re.I
        )

        def _shape_type_str(sh):
            try:
                return str(sh.shape_type)
            except Exception:
                return ""

        def _rotation_deg(sh):
            try:
                rot = getattr(sh, "rotation", None)
                if rot is None:
                    return 0.0
                return float(rot) % 360.0
            except Exception:
                return 0.0

        def _preset_geom(sh):
            """取 shape 的 a:prstGeom prst 值（若存在）。"""
            try:
                el = sh._element.find(".//a:prstGeom", NS)
                if el is not None:
                    return (el.get("prst") or "").lower()
            except Exception:
                pass
            return ""

        def _looks_like_pointer(sh):
            """形状语义线索：三角/箭头/名称含关键词。"""
            st = _shape_type_str(sh)
            prst = _preset_geom(sh)
            name = sh.name or ""
            reasons = []
            # PowerPoint 常见 prst：triangle / rtTriangle / downArrow / upArrow /
            # leftArrow / rightArrow / arrowCallout* 等
            if prst:
                if "triangle" in prst or "arrow" in prst:
                    reasons.append(f"prst={prst}")
            if POINTER_SHAPE_HINT.search(name):
                reasons.append(f"name={name!r}")
            if "TRIANGLE" in st or "ARROW" in st or "FREEFORM" in st:
                # freeform 不足以单独作为指针证据，但可与朝向证据组合
                if "TRIANGLE" in st or "ARROW" in st:
                    reasons.append(f"shape_type={st}")
                else:
                    reasons.append(f"shape_type={st}(需与朝向组合)")
            return reasons

        def _tip_xy(sh):
            """估算指针尖端在幻灯片上的 (x_cm, y_cm)。
            规则：
              · 起始参考："默认尖端"在 bbox 顶边中点（即"未旋转的向上箭头"，
                尖端朝上；因为指针位于转盘上方，尖端应朝下指向转盘顶边）。
                因此我们改用"底边中点"作为向下箭头尖端的默认位置。
              · 若形状 prst 明确了朝向（如 downArrow），则底边中点；如 upArrow
                则顶边中点；其余根据 rotation 旋转 bbox 内的默认位置。
              · rotation 用 shape 的 rotation 属性；正值为顺时针（度）。
            返回 (tip_x_cm, tip_y_cm, orient_desc)。
            """
            left = emu_to_cm(sh.left or 0)
            top = emu_to_cm(sh.top or 0)
            w, h = shape_size_cm(sh)
            cx = left + w / 2.0
            cy = top + h / 2.0
            prst = _preset_geom(sh)
            rot = _rotation_deg(sh)
            # 默认"尖端方向向量"（bbox 局部坐标系，向下 = (0, +h/2)）
            # 之所以默认向下：转盘位于指针下方，正确的指针尖端应朝下。
            local_dx, local_dy = 0.0, h / 2.0
            orient = "default-down"
            if "uparrow" in prst or "uptriangle" in prst or prst == "triangle":
                # PPT 的 "triangle" prst 是尖端朝上的等腰三角形
                local_dx, local_dy = 0.0, -h / 2.0
                orient = f"prst={prst}(up)"
            elif "downarrow" in prst or "downtriangle" in prst:
                local_dx, local_dy = 0.0, h / 2.0
                orient = f"prst={prst}(down)"
            elif "leftarrow" in prst:
                local_dx, local_dy = -w / 2.0, 0.0
                orient = f"prst={prst}(left)"
            elif "rightarrow" in prst:
                local_dx, local_dy = w / 2.0, 0.0
                orient = f"prst={prst}(right)"
            # 应用旋转（顺时针为正）
            import math
            theta = math.radians(rot)
            rx = local_dx * math.cos(theta) - local_dy * math.sin(theta)
            ry = local_dx * math.sin(theta) + local_dy * math.cos(theta)
            return cx + rx, cy + ry, f"{orient} rot={rot:.0f}°"

        def _classify_pointer_for(spinner, radius, sp_center, all_cands):
            """在候选中找一个既位于转盘上方，又"指向轮盘边缘"的指针。"""
            sx, sy = sp_center
            # 转盘顶边 y = sy - radius
            top_edge_y = sy - radius
            best = None
            best_debug = ""
            for s in all_cands:
                if s in spinners:
                    continue
                cx, cy = shape_center_cm(s)
                w, h = shape_size_cm(s)
                # 粗筛：中心在转盘上方（cy 小于转盘中心）
                if cy >= sy:
                    continue
                # 粗筛：x 大致在转盘 x 范围（±radius+1cm）
                if abs(cx - sx) > radius + 1.0:
                    continue
                # 形状语义证据
                reasons = _looks_like_pointer(s)
                if not reasons:
                    continue
                # 朝向：bbox 细长比 或 明确 prst
                aspect_ok = max(w, h) / max(min(w, h), 1e-6) >= 1.4
                prst = _preset_geom(s)
                has_prst_dir = any(k in prst for k in
                                   ("uparrow", "downarrow", "leftarrow",
                                    "rightarrow", "triangle"))
                if not (aspect_ok or has_prst_dir):
                    continue
                # 尖端位置：需要落在转盘顶边附近
                tip_x, tip_y, orient = _tip_xy(s)
                y_tol = max(radius * 0.35, 0.5)
                x_tol = max(radius * 0.35, 0.5)
                y_ok = abs(tip_y - top_edge_y) <= y_tol
                x_ok = abs(tip_x - sx) <= x_tol
                dbg = (f"name={s.name!r} prst={prst!r} 中心=({cx:.2f},{cy:.2f}) "
                       f"tip=({tip_x:.2f},{tip_y:.2f}) 转盘顶=({sx:.2f},{top_edge_y:.2f}) "
                       f"orient={orient} 语义={reasons} y_ok={y_ok} x_ok={x_ok}")
                if y_ok and x_ok:
                    best = s
                    best_debug = dbg
                    break
                else:
                    if not best_debug:
                        best_debug = "[未通过] " + dbg
            return best, best_debug

        left_ptr, left_dbg = _classify_pointer_for(left_sp, l_rad, cl_sp, pointer_cands)
        right_ptr, right_dbg = _classify_pointer_for(right_sp, r_rad, cr_sp, pointer_cands)
        rule.notes.append(f"左转盘指针：{left_dbg or '(无候选通过)'}")
        rule.notes.append(f"右转盘指针：{right_dbg or '(无候选通过)'}")

        if left_ptr is None or right_ptr is None:
            rule.notes.append("未同时找到两个通过'形状/朝向/尖端在轮盘边缘'检验的指针")
            self.report.rules.append(rule)
            return

        # ---- 大小一致（宽高差 ≤0.2cm） ---- #
        wL, hL = shape_size_cm(left_ptr)
        wR, hR = shape_size_cm(right_ptr)
        size_ok = abs(wL - wR) <= 0.2 and abs(hL - hR) <= 0.2

        # ---- 位置一致：相对各自转盘中心的偏移量接近（容差 0.3cm） ---- #
        cL = shape_center_cm(left_ptr)
        cR = shape_center_cm(right_ptr)
        off_xL = cL[0] - cl_sp[0]
        off_yL = cL[1] - cl_sp[1]
        off_xR = cR[0] - cr_sp[0]
        off_yR = cR[1] - cr_sp[1]
        pos_ok = abs(off_xL - off_xR) <= 0.3 and abs(off_yL - off_yR) <= 0.3

        rule.notes.append(
            f"左指针 wh=({wL:.2f},{hL:.2f}) 偏移=({off_xL:.2f},{off_yL:.2f}); "
            + f"右指针 wh=({wR:.2f},{hR:.2f}) 偏移=({off_xR:.2f},{off_yR:.2f}); "
            + f"大小一致={size_ok} 位置一致={pos_ok}"
        )
        rule.matched = size_ok and pos_ok
        self.report.rules.append(rule)

    def _rule_cross_format(self, info):
        """+1 可单独选中拖动的 × 对象，字号约80~90磅，加粗，字体颜色白/米/淡粉。

        严格化（回应反馈）：
          - 字号 / 加粗 / 颜色：从 run 属性开始，逐层回退到段落 defRPr、
            txBody 的 lstStyle、占位符继承（layout/master）、以及 theme
            的字体默认，直到取到有效值；颜色支持 srgbClr / schemeClr /
            prstClr 三种，并可通过 lumMod/lumOff 亮度修正。
          - 可单独选中拖动：显式检查 shape 上是否含 spLocks/grpSpLocks 的
            noSelect / noMove / noGrp 等锁定属性；命中即视为不可单独拖动。
        """
        rule = Rule("+1-cross-format", 1,
                    "× 可单独选中(未锁定)，字号80~90磅，加粗，"
                    + "颜色白/米/淡粉（含 theme/继承颜色解析）")

        if info["fullpage_media"] and not info["cross_x"]:
            rule.notes.append("存在整页视频且页面无独立 × 对象，跳过不计分")
            self.report.rules.append(rule)
            return

        if not info["cross_x"]:
            rule.notes.append("未发现 × 文本对象")
            self.report.rules.append(rule)
            return

        for s in info["cross_x"]:
            # ---- 锁定检查（可单独选中拖动） ---- #
            try:
                sp_xml = ET.tostring(s._element, encoding="unicode")
            except Exception:
                sp_xml = ""
            lock_hit = None
            for attr in ('noSelect="1"', 'noSelect="true"',
                         'noMove="1"', 'noMove="true"',
                         'noGrp="1"', 'noGrp="true"'):
                if attr in sp_xml:
                    lock_hit = attr
                    break
            if lock_hit:
                rule.notes.append(
                    f"× 对象 name={s.name!r} 存在锁定属性 {lock_hit}，不满足'可单独选中拖动'"
                )
                continue

            # ---- 逐 run 解析有效字号/加粗/颜色（含继承与 theme） ---- #
            try:
                tf = s.text_frame
            except Exception:
                continue
            hit_any = False
            fail_details = []
            for para in tf.paragraphs:
                for run in para.runs:
                    eff = self._resolve_run_effective_style(s, para, run)
                    sz = eff["size_pt"]
                    bold = eff["bold"]
                    color_hex = eff["color_hex"]
                    reason = eff["reason"]
                    ok_sz = sz is not None and 80 <= sz <= 90
                    ok_bold = bool(bold)
                    ok_color = color_hex is not None and self._is_white_beige_pink(color_hex)
                    if ok_sz and ok_bold and ok_color:
                        rule.notes.append(
                            f"× name={s.name!r} 命中：字号={sz}pt 粗体={bold} "
                            + f"颜色=#{color_hex}  解析路径={reason}"
                        )
                        hit_any = True
                        break
                    fail_details.append(
                        f"字号={sz} 粗体={bold} 颜色={color_hex} "
                        + f"[sz_ok={ok_sz} bold_ok={ok_bold} color_ok={ok_color}] "
                        + f"路径={reason}"
                    )
                if hit_any:
                    break
            if hit_any:
                rule.matched = True
                self.report.rules.append(rule)
                return
            rule.notes.append(f"× name={s.name!r} 未命中：" + " | ".join(fail_details))
        self.report.rules.append(rule)

    # ---------- 文本样式与主题解析 ---------- #

    def _resolve_run_effective_style(self, shape, paragraph, run) -> dict:
        """按 OOXML 继承链解析一个 run 的有效 字号 / 加粗 / 颜色。

        继承顺序（前者优先）：
          1. run 直接属性（a:rPr）
          2. 段落 defRPr（a:pPr/a:defRPr）
          3. shape 的 txBody/lstStyle 中对应级别（a:lvlNPpr/a:defRPr，
             以及 a:defPPr/a:defRPr）
          4. 占位符：slide layout 的对应占位符 -> master 的对应占位符
          5. 主题默认（theme/fontScheme + clrScheme）

        颜色支持：a:srgbClr / a:schemeClr（通过 theme 的 clrScheme 解析）/
                   a:prstClr；带 a:lumMod / a:lumOff 亮度修正。

        返回：{"size_pt": float|None, "bold": bool|None, "color_hex": str|None,
                "reason": str（解析出处，便于 debug）}
        """
        reason_parts = []
        # 快照：从 run 开始，逐级找 rPr
        run_rpr = self._safe_find(getattr(run, "_r", None), "a:rPr")
        para_el = getattr(paragraph, "_p", None)
        para_pPr = self._safe_find(para_el, "a:pPr") if para_el is not None else None
        para_defrpr = self._safe_find(para_pPr, "a:defRPr") if para_pPr is not None else None

        size_pt = None
        bold = None
        color_hex = None

        def _apply_rpr(rpr, tag):
            nonlocal size_pt, bold, color_hex
            if rpr is None:
                return
            # size
            if size_pt is None:
                sz_attr = rpr.get("sz")
                if sz_attr and sz_attr.isdigit():
                    size_pt = int(sz_attr) / 100.0
                    reason_parts.append(f"size<-{tag}.sz={sz_attr}")
            # bold
            if bold is None:
                b_attr = rpr.get("b")
                if b_attr is not None:
                    bold = b_attr in ("1", "true", "True")
                    reason_parts.append(f"bold<-{tag}.b={b_attr}")
            # color
            if color_hex is None:
                fill_hex = self._extract_solid_fill_hex(rpr)
                if fill_hex:
                    color_hex = fill_hex.upper()
                    reason_parts.append(f"color<-{tag}.solidFill={color_hex}")

        # 1) run
        _apply_rpr(run_rpr, "run.rPr")
        # 2) 段落 defRPr
        _apply_rpr(para_defrpr, "para.defRPr")

        # 3) shape txBody 的 lstStyle
        try:
            tx_body = shape._element.find(".//p:txBody", NS)
        except Exception:
            tx_body = None
        if tx_body is not None:
            lst_style = tx_body.find("a:lstStyle", NS)
            if lst_style is not None:
                # 取段落级：优先 lvl1pPr，其次 defPPr
                for lvl_tag in ("a:lvl1pPr", "a:defPPr", "a:lvl2pPr",
                                "a:lvl3pPr", "a:lvl4pPr", "a:lvl5pPr"):
                    lvl = lst_style.find(lvl_tag, NS)
                    if lvl is not None:
                        _apply_rpr(lvl.find("a:defRPr", NS), f"lstStyle.{lvl_tag}.defRPr")

        # 4) 占位符继承（layout/master）
        if size_pt is None or bold is None or color_hex is None:
            ph_rpr_list = self._collect_placeholder_defrpr(shape)
            for src_tag, rpr in ph_rpr_list:
                _apply_rpr(rpr, src_tag)
                if size_pt is not None and bold is not None and color_hex is not None:
                    break

        # 5) 主题默认（bodyPr/majorFont/minorFont + clrScheme dk1）
        if color_hex is None:
            # 若颜色仍为 None，取 theme 的 dk1（正文默认色）
            theme_color = self._theme_scheme_color("dk1")
            if theme_color:
                color_hex = theme_color.upper()
                reason_parts.append("color<-theme.dk1")

        return {
            "size_pt": size_pt,
            "bold": bold,
            "color_hex": color_hex,
            "reason": " > ".join(reason_parts) or "(默认)",
        }

    def _safe_find(self, el, tag):
        try:
            if el is None:
                return None
            return el.find(tag, NS)
        except Exception:
            return None

    def _extract_solid_fill_hex(self, rpr) -> Optional[str]:
        """从 rPr 里解析 a:solidFill 的颜色为 6 位 RGB hex（考虑 lumMod/lumOff）。"""
        if rpr is None:
            return None
        sf = rpr.find("a:solidFill", NS)
        if sf is None:
            return None
        # srgbClr
        srgb = sf.find("a:srgbClr", NS)
        if srgb is not None:
            hexv = (srgb.get("val") or "").upper()
            if len(hexv) == 6:
                return self._apply_lum_mod_off(hexv, srgb)
        # schemeClr
        sc = sf.find("a:schemeClr", NS)
        if sc is not None:
            name = sc.get("val") or ""
            resolved = self._theme_scheme_color(name)
            if resolved:
                return self._apply_lum_mod_off(resolved.upper(), sc)
        # prstClr
        pr = sf.find("a:prstClr", NS)
        if pr is not None:
            hexv = self._preset_color_hex(pr.get("val") or "")
            if hexv:
                return self._apply_lum_mod_off(hexv, pr)
        return None

    @staticmethod
    def _apply_lum_mod_off(hexv: str, color_el) -> str:
        """对 hex 颜色套用 <a:lumMod val="..."/> 和 <a:lumOff val="..."/> 亮度修正。
        PowerPoint OOXML 中 val 为千分数（60000 = 60%）。
        """
        try:
            r = int(hexv[0:2], 16) / 255.0
            g = int(hexv[2:4], 16) / 255.0
            b = int(hexv[4:6], 16) / 255.0
        except Exception:
            return hexv
        lum_mod = color_el.find("a:lumMod", NS)
        lum_off = color_el.find("a:lumOff", NS)
        # 用 HSL 亮度近似：L = (max+min)/2
        mx = max(r, g, b)
        mn = min(r, g, b)
        L = (mx + mn) / 2.0
        new_L = L
        if lum_mod is not None:
            try:
                new_L = new_L * (int(lum_mod.get("val") or "100000") / 100000.0)
            except Exception:
                pass
        if lum_off is not None:
            try:
                new_L = new_L + (int(lum_off.get("val") or "0") / 100000.0)
            except Exception:
                pass
        new_L = max(0.0, min(1.0, new_L))
        # 缩放 RGB：粗略保持色相 (r,g,b) 按新 L 与旧 L 比例调整
        if L <= 0:
            r2 = g2 = b2 = new_L
        else:
            scale = new_L / L if L > 0 else 1.0
            r2 = min(1.0, r * scale)
            g2 = min(1.0, g * scale)
            b2 = min(1.0, b * scale)
        return "{:02X}{:02X}{:02X}".format(int(r2 * 255), int(g2 * 255), int(b2 * 255))

    def _collect_placeholder_defrpr(self, shape):
        """若 shape 是占位符，从 slide layout / master 上找到对应占位符的
        默认段落 defRPr，返回 [(source_tag, rpr_element), ...]。
        """
        out = []
        try:
            ph = shape._element.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
        except Exception:
            ph = None
        if ph is None:
            return out
        ph_type = ph.get("type") or ""
        ph_idx = ph.get("idx") or ""

        # slide layout
        try:
            layout = self.slide.slide_layout
        except Exception:
            layout = None
        for src_name, src in (("layout", layout),
                              ("master", layout.slide_master if layout is not None else None)):
            if src is None:
                continue
            try:
                spTree = src.element.find(".//p:cSld/p:spTree", NS)
            except Exception:
                spTree = None
            if spTree is None:
                continue
            for sp in spTree.findall("p:sp", NS):
                sp_ph = sp.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
                if sp_ph is None:
                    continue
                if ph_type and sp_ph.get("type") == ph_type:
                    match = True
                elif ph_idx and sp_ph.get("idx") == ph_idx:
                    match = True
                else:
                    match = False
                if not match:
                    continue
                lst = sp.find(".//p:txBody/a:lstStyle", NS)
                if lst is not None:
                    for lvl_tag in ("a:lvl1pPr", "a:defPPr"):
                        lvl = lst.find(lvl_tag, NS)
                        if lvl is not None:
                            rpr = lvl.find("a:defRPr", NS)
                            if rpr is not None:
                                out.append((f"{src_name}.ph.{lvl_tag}.defRPr", rpr))
        return out

    def _theme_scheme_color(self, name: str) -> Optional[str]:
        """从演示文稿主题的 clrScheme 中解析 scheme color name（dk1/lt1/
        accent1..6/hlink/folHlink）为 6 位 hex。缓存在 self 上。
        """
        if not name:
            return None
        cache = getattr(self, "_theme_cache", None)
        if cache is None:
            cache = self._load_theme_scheme()
            self._theme_cache = cache
        # 兼容 dk1/tx1、lt1/bg1 等别名
        aliases = {
            "tx1": "dk1", "bg1": "lt1",
            "tx2": "dk2", "bg2": "lt2",
        }
        key = aliases.get(name, name)
        return cache.get(key)

    def _load_theme_scheme(self) -> dict:
        """加载 ppt/theme/theme1.xml，返回 {schemeName: hex6}。"""
        out = {}
        if self.zf is None:
            return out
        theme_name = None
        for n in self.zf.namelist():
            if n.startswith("ppt/theme/") and n.endswith(".xml"):
                theme_name = n
                break
        if theme_name is None:
            return out
        try:
            root = ET.fromstring(self.zf.read(theme_name))
        except Exception:
            return out
        scheme = root.find(".//a:themeElements/a:clrScheme", NS)
        if scheme is None:
            return out
        for child in scheme:
            tag = child.tag.split("}")[-1]  # dk1, lt1, accent1, ...
            srgb = child.find("a:srgbClr", NS)
            if srgb is not None:
                v = (srgb.get("val") or "").upper()
                if len(v) == 6:
                    out[tag] = v
                    continue
            # 系统颜色（sysClr）：使用 lastClr 属性
            sysc = child.find("a:sysClr", NS)
            if sysc is not None:
                v = (sysc.get("lastClr") or "").upper()
                if len(v) == 6:
                    out[tag] = v
        return out

    @staticmethod
    def _preset_color_hex(name: str) -> Optional[str]:
        """OOXML a:prstClr 的常见预设颜色到 hex 的映射（部分常用值）。"""
        table = {
            "white": "FFFFFF", "black": "000000",
            "beige": "F5F5DC", "bisque": "FFE4C4",
            "mistyRose": "FFE4E1", "lightPink": "FFB6C1",
            "pink": "FFC0CB", "hotPink": "FF69B4",
            "wheat": "F5DEB3", "cornsilk": "FFF8DC",
            "ivory": "FFFFF0", "linen": "FAF0E6",
            "seashell": "FFF5EE", "papayaWhip": "FFEFD5",
            "peachPuff": "FFDAB9", "antiqueWhite": "FAEBD7",
        }
        return table.get(name)

    @staticmethod
    def _is_white_beige_pink(hexstr: str) -> bool:
        try:
            r = int(hexstr[0:2], 16)
            g = int(hexstr[2:4], 16)
            b = int(hexstr[4:6], 16)
        except Exception:
            return False
        # 白色
        if r >= 240 and g >= 240 and b >= 240:
            return True
        # 米色 (Beige ~ F5F5DC, 偏黄白)
        if r >= 220 and g >= 210 and 180 <= b <= 230 and r >= b:
            return True
        # 淡粉
        if r >= 230 and 180 <= g <= 230 and 200 <= b <= 240 and r >= g:
            return True
        return False

    def _rule_button_states(self, info):
        """+1 按钮根据状态显示 开始/暂停/继续 或语义一致提示。

        严格化（回应反馈）：rubric 只要求"根据状态显示"三类语义，
        单按钮初始只显示"开始"、点击后动态变更也是合格的。
        因此判定顺序：
          (A) 若存在多按钮/多状态文本，直接根据当前静态文本命中：
                只要 {开始|暂停|继续} 三类语义中至少两类有静态可见证据，即命中。
          (B) 单按钮或静态只见一类状态词时，尝试从动画/宏中证明状态序列：
              · p:timing 中该按钮 onClick 触发的 par 数 ≥ 2（点击 → 停 → 继）
              · 或 par 内有 pause/resume/togglePause 之一，且另一 par 有反向命令
              · 或按钮绑定 VBA(ppaction://macro / vbaProject.bin) —— 无法静态
                展开脚本，但至少能证明"点击后状态会变化"；给以中立命中。
          (C) 无任何证据时保守判为不命中。
        """
        rule = Rule("+1-button-states", 1,
                    "按钮状态文字显示开始/暂停/继续或语义一致提示"
                    + "（含动态状态/序列证据）")

        # 整页视频存在且无独立按钮，跳过不计分
        if info["fullpage_media"] and not info["buttons"]:
            rule.notes.append("存在整页视频且页面无独立按钮，跳过不计分")
            self.report.rules.append(rule)
            return

        if not info["buttons"]:
            rule.notes.append("未发现按钮候选")
            self.report.rules.append(rule)
            return

        start_words = ("开始", "播放", "启动", "start", "play", "begin", "go")
        pause_words = ("暂停", "停止", "pause", "stop", "halt")
        resume_words = ("继续", "恢复", "resume", "continue")

        def _cls_of(text: str) -> "set[str]":
            t = (text or "").lower()
            cls = set()
            if any(w.lower() in t for w in start_words):
                cls.add("start")
            if any(w.lower() in t for w in pause_words):
                cls.add("pause")
            if any(w.lower() in t for w in resume_words):
                cls.add("resume")
            return cls

        # ---- (A) 静态文本命中 ---- #
        classes_seen: "set[str]" = set()
        per_button = []
        for s in info["buttons"]:
            text = safe_text(s)
            cls = _cls_of(text)
            classes_seen |= cls
            per_button.append((s, text, cls))
        rule.notes.append(
            "按钮静态文本："
            + " | ".join(f"{safe_text(s).strip()!r}->{sorted(c)}"
                          for s, _t, c in per_button)
        )

        # 若已见 ≥2 类静态语义 -> 直接命中
        if len(classes_seen) >= 2:
            rule.notes.append(f"静态文本已覆盖状态语义 {sorted(classes_seen)}，命中")
            rule.matched = True
            self.report.rules.append(rule)
            return

        # ---- (B) 从动画/VBA 证明状态序列 ---- #
        # 逐个可能是"控制型"的按钮尝试匹配
        candidates = []
        for s, text, cls in per_button:
            # 只对像"控制按钮"的候选做动态分析：文本已含至少一类状态词，
            # 或名称含 button/btn/控制/按钮
            name = (s.name or "").lower()
            if cls or re.search(r"(button|btn|控制|按钮)", name):
                candidates.append(s)
        if not candidates:
            candidates = [s for s, _t, _c in per_button]

        dynamic_hit = False
        dynamic_reason = None
        for btn in candidates:
            # 静态无法解析 VBA 时给中立命中（有 vbaProject.bin + program 动作）
            try:
                btn_xml = ET.tostring(btn._element, encoding="unicode")
            except Exception:
                btn_xml = ""
            macro_action = bool(re.search(r'action="ppaction://macro[^"]*"', btn_xml))
            has_vba = False
            try:
                if self.zf is not None and "ppt/vbaProject.bin" in self.zf.namelist():
                    has_vba = True
            except Exception:
                pass
            program_action = bool(re.search(r'action="ppaction://program', btn_xml))
            if macro_action or (has_vba and program_action):
                dynamic_hit = True
                dynamic_reason = (
                    f"按钮 name={btn.name!r} 绑定 VBA/自定义程序动作 "
                    + "-> 状态文本会由脚本动态更新（静态无法展开，判为合格）"
                )
                break

            # 分析 timing：该按钮 onClick 触发的 par 数 ≥ 2 或含 pause/resume 命令
            if self.slide_xml is None:
                continue
            timing = self.slide_xml.find(".//p:timing", NS)
            if timing is None:
                continue
            btn_spid = None
            for tag in ("p:nvSpPr/p:cNvPr", "p:nvPicPr/p:cNvPr", "p:nvGrpSpPr/p:cNvPr"):
                el = btn._element.find(".//" + tag, NS)
                if el is not None:
                    btn_spid = el.get("id")
                    break
            if not btn_spid:
                continue

            click_pars = []
            for par in timing.iter("{%s}par" % NS["p"]):
                for cond in par.findall(".//p:cond", NS):
                    if cond.get("evt") != "onClick":
                        continue
                    sptgt = cond.find(".//p:spTgt", NS)
                    if sptgt is not None and sptgt.get("spid") == btn_spid:
                        click_pars.append(par)
                        break

            has_pause_cmd = False
            has_resume_cmd = False
            for par in click_pars:
                for cmd in par.findall(".//p:cmd", NS):
                    ctype = (cmd.get("type") or "").lower()
                    if ctype in ("stop", "pause"):
                        has_pause_cmd = True
                    if ctype in ("play", "resume", "togglepause"):
                        has_resume_cmd = True
                if par.find(".//p:cmd[@type='togglePause']", NS) is not None:
                    has_pause_cmd = True
                    has_resume_cmd = True

            if len(click_pars) >= 2 or (has_pause_cmd and has_resume_cmd):
                dynamic_hit = True
                dynamic_reason = (
                    f"按钮 name={btn.name!r} 有 {len(click_pars)} 个 onClick par "
                    + f"(pause_cmd={has_pause_cmd}, resume_cmd={has_resume_cmd}) "
                    + "-> 可证明按钮点击后状态会切换，视为动态状态文本"
                )
                break
            if len(click_pars) >= 1 and classes_seen:
                # 单 par 但静态已见 ≥1 类状态词：认为初始 UI 显示"开始"、
                # 点击后由动画驱动的隐藏/显示切换成"暂停/继续"，也认合格。
                dynamic_hit = True
                dynamic_reason = (
                    f"按钮 name={btn.name!r} 静态见 {sorted(classes_seen)}，"
                    + f"且有 {len(click_pars)} 个 onClick par -> 判为动态状态切换"
                )
                break

        if dynamic_hit:
            rule.notes.append(dynamic_reason or "动态状态证据命中")
            rule.matched = True
            self.report.rules.append(rule)
            return

        # ---- (C) 全部证据缺失 ---- #
        rule.notes.append(
            "无静态多语义文本，也无动画/VBA 状态切换证据 -> 判为不命中"
        )
        self.report.rules.append(rule)

    # ---------- 扣分项实现 ---------- #

    def _rule_minus_occlusion(self, info, sw, sh):
        """-3 新增对象遮挡粉色椭圆、乘号、按钮或页面主要装饰元素"""
        rule = Rule("-3-occlusion", -3, "新增对象遮挡椭圆/乘号/按钮/装饰")

        if self.slide is None:
            self.report.rules.append(rule)
            return

        # 整页视频存在且没有其它可核验遮挡对象，跳过不计分
        non_fullpage_shapes = [
            s for s in self.slide.shapes
            if s not in info["fullpage_media"] and s not in info["fullpage_picture"]
        ]
        if info["fullpage_media"] and not non_fullpage_shapes:
            rule.notes.append("存在整页视频且无其它对象，跳过不计分")
            self.report.rules.append(rule)
            return

        # 只针对细则：新增对象遮挡粉色椭圆、乘号、按钮或页面主要装饰元素。
        # 自动化代理：顶层非目标对象若大面积覆盖页面中央区域，视为遮挡主要元素。
        if non_fullpage_shapes:
            last = non_fullpage_shapes[-1]
            w, h = shape_size_cm(last)
            cx, cy = shape_center_cm(last)
            covers_center = (sw * 0.25 <= cx <= sw * 0.75) and (sh * 0.25 <= cy <= sh * 0.75)
            large_object = w >= sw * 0.5 or h >= sh * 0.5
            target_like = last in info["cross_x"] or last in info["buttons"]
            if large_object and covers_center and not target_like:
                rule.notes.append(f"顶层新增对象 '{last.name}' 覆盖页面主要区域")
                rule.matched = True

        if not rule.matched:
            rule.notes.append("未发现新增对象遮挡指定元素")
        self.report.rules.append(rule)

    def _rule_minus_fullpage_video(self, info):
        """-5 第10页被整页视频或动画覆盖，或转盘/按钮无法单独编辑"""
        rule = Rule("-5-fullpage-video", -5, "整页视频/动画覆盖，无法单独编辑")

        # 条件 1：出现整页视频或动画覆盖
        if info["fullpage_media"]:
            rule.notes.append(f"发现整页媒体：{info['fullpage_media'][0].name}")
            rule.matched = True

        # 条件 2：没有整页视频时，转盘或按钮对象被锁定，无法单独编辑
        def is_locked(shape) -> bool:
            xml = ET.tostring(shape._element, encoding="unicode")
            lock_attrs = (
                'noSelect="1"', 'noSelect="true"',
                'noMove="1"', 'noMove="true"',
                'noResize="1"', 'noResize="true"',
                'noEditPoints="1"', 'noEditPoints="true"',
                'noGrp="1"', 'noGrp="true"',
                'noChangeAspect="1"', 'noChangeAspect="true"',
            )
            return any(attr in xml for attr in lock_attrs)

        spinner_candidates = []
        for s in (info["ovals"] + info["pictures"] + info["groups"]):
            if s in info["fullpage_media"] or s in info["fullpage_picture"]:
                continue
            w, h = shape_size_cm(s)
            if w > 0 and h > 0 and abs(w - h) <= 0.2:
                spinner_candidates.append(s)

        locked_spinner = any(is_locked(s) for s in spinner_candidates)
        locked_button = any(is_locked(s) for s in info["buttons"])
        if locked_spinner or locked_button:
            rule.notes.append(
                f"检测到无法单独编辑对象：转盘锁定={locked_spinner}，按钮锁定={locked_button}"
            )
            rule.matched = True

        if not rule.matched:
            rule.notes.append("未发现整页视频/动画，未发现转盘或按钮锁定")
        self.report.rules.append(rule)

    def _rule_minus_garbage(self, info):
        """-1 文件中出现批注/临时说明/红色标记/截图边框/无关占位对象"""
        rule = Rule("-1-garbage", -1, "批注/临时说明/红色标记/截图边框/无关占位")
        # 1) 批注：演讲者备注（notesSlide10.xml）中是否含实质内容
        notes_text = ""
        try:
            notes_xml = self.zf.read(
                f"ppt/notesSlides/notesSlide{self.TARGET_PAGE_INDEX + 1}.xml"
            ).decode("utf-8")
            # 提取 <a:t> 文本
            notes_text = " ".join(re.findall(r"<a:t>(.*?)</a:t>", notes_xml))
        except KeyError:
            pass
        meaningful_notes = bool(notes_text.strip()) and len(notes_text.strip()) > 5
        rule.notes.append(f"备注页文本长度={len(notes_text.strip())}: {notes_text[:80]!r}")
        # 2) PPT 批注（comments）
        has_comments = False
        for n in self.zf.namelist():
            if n.startswith("ppt/comments/"):
                has_comments = True
                break
        rule.notes.append(f"含 PPT 批注：{has_comments}")
        # 3) 红色标记：检测页面 shape 中是否存在显著红色填充/线条小图形（颜色 R>200, G<80, B<80）
        red_found = False
        red_pat = re.compile(r'srgbClr val="([0-9A-Fa-f]{6})"')
        for m in red_pat.finditer(self.slide_xml_raw):
            r = int(m.group(1)[0:2], 16)
            g = int(m.group(1)[2:4], 16)
            b = int(m.group(1)[4:6], 16)
            if r > 200 and g < 80 and b < 80:
                red_found = True
                break
        rule.notes.append(f"含纯红色 srgb 标记：{red_found}")
        # 4) 截图边框：图片名 含 screenshot/截图，或图片是带白边的矩形（这里只判文件名）
        screenshot_like = False
        for s in info["pictures"]:
            if re.search(r"screenshot|截图|snip", s.name or "", re.I):
                screenshot_like = True
                break
        rule.notes.append(f"截图类图片：{screenshot_like}")
        # 5) 无关占位：layout 占位符未删除（在 slide xml 中有 <p:ph>）
        ph_count = len(re.findall(r"<p:ph\b", self.slide_xml_raw))
        rule.notes.append(f"页面占位符数：{ph_count}")

        rule.matched = (
            meaningful_notes or has_comments or red_found or screenshot_like or ph_count > 0
        )
        self.report.rules.append(rule)

    # -------- 入口 -------- #

    def run(self) -> Report:
        if not self.check_dim1():
            return self.report
        self.report.dim1_pass = True
        self.check_dim2()
        return self.report


# ----------------------------- 统一入口 ----------------------------- #


SCRIPT_ID = "065"


def _locate_document(dir_path: str) -> Optional[str]:
    """在给定目录中定位待评估的 .pptx 文档；找不到返回 None。
    忽略 Office 打开时生成的临时锁文件（以 ~$ 开头）。
    仅支持 .pptx；旧版 .ppt 不再识别（避免依赖 COM/LibreOffice 转换）。
    """
    try:
        entries = os.listdir(dir_path)
    except OSError:
        return None
    for name in entries:
        if name.startswith("~$"):
            continue
        if name.lower().endswith(".pptx"):
            return os.path.join(dir_path, name)
    return None


def evaluate(dir_path: str) -> dict:
    """统一入口：接收脚本所在目录的路径，脚本自行在该目录中定位并评估文档。

    返回结构见《脚本接口差异与统一建议.md》§2.2。
    命中和未命中的维度二细则都会返回，便于批量汇总。
    """
    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": 0,
    }

    try:
        if not os.path.isdir(dir_path):
            result["status"] = "error"
            result["error"] = f"目录不存在：{dir_path}"
            return result

        file_path = _locate_document(dir_path)
        if file_path is None:
            result["status"] = "error"
            result["error"] = f"目录中未找到 .pptx 文件：{dir_path}"
            return result
        result["file_name"] = os.path.basename(file_path)

        rep = Evaluator(file_path).run()

        result["dim1_pass"] = rep.dim1_pass
        if not rep.dim1_pass:
            # 维度一未通过：把关键原因串起来供人工回看
            result["dim1_reason"] = "；".join(rep.dim1_reasons)

        items = []
        total = 0
        max_score = 0
        for r in rep.rules:
            delta = r.score if r.matched else 0
            items.append({
                "rule": r.title,
                "max_delta": r.score,
                "delta": delta,
                "hit": r.matched,
                "detail": "",
            })
            total += delta
            if r.score > 0:
                max_score += r.score

        result["dim2_items"] = items
        result["total_score"] = total
        result["max_score"] = max_score
    except Exception as e:
        result["status"] = "error"
        result["error"] = f"{type(e).__name__}: {e}"

    return result


if __name__ == "__main__":
    # 仅本地调试：传入脚本所在目录（默认取脚本自身所在目录）
    _dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(
        os.path.abspath(__file__)
    )
    print(json.dumps(evaluate(_dir), ensure_ascii=False, indent=2))
