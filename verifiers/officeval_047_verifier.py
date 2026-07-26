#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# pyright: reportUnknownMemberType=false, reportUnknownVariableType=false, reportUnknownArgumentType=false, reportUnknownParameterType=false, reportMissingTypeArgument=false, reportImplicitStringConcatenation=false, reportDeprecated=false, reportExplicitAny=false, reportGeneralTypeIssues=false, reportArgumentType=false, reportAttributeAccessIssue=false, reportPrivateUsage=false, reportOptionalMemberAccess=false, reportOptionalSubscript=false, reportOptionalIterable=false, reportOptionalCall=false
"""自动评估《设备状态监测研究汇报》PPT。

统一接口约定（详见《脚本接口差异与统一建议》§2）：
    · 对外仅暴露 evaluate(dir_path: str) -> dict
    · 参数为脚本所在目录的路径，脚本自行在该目录中定位待评估文档
    · 返回结构化字典（含维度一通过与否、维度二逐项得分、总分）
    · 主结果只走 return，不 print、不改 sys.stdout、不 sys.exit、不硬编码路径

评分口径不变：维度一有任一项未通过 → 直接 0 分且不再评维度二；
维度二按得分点与扣分点累加得到 total_score。
"""

from __future__ import annotations

import json
import math
import os
import re
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

EMU_PER_CM = 360_000
PT_PER_EMU = 12700
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}
SANS_FONTS = (
    "微软雅黑",
    "microsoft yahei",
    "黑体",
    "simhei",
    "arial",
    "calibri",
    "等线",
    "dengxian",
    "汉仪元隆黑",
    "hy yuanlong",
)


@dataclass
class Box:
    left: float
    top: float
    width: float
    height: float

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height

    @property
    def cx(self) -> float:
        return self.left + self.width / 2

    @property
    def cy(self) -> float:
        return self.top + self.height / 2

    @property
    def area(self) -> float:
        return max(0.0, self.width) * max(0.0, self.height)


@dataclass
class RuleResult:
    name: str
    score: int
    hit: bool
    detail: str


class PPTInspector:
    def __init__(self, path: Path):
        self.path = path
        self.prs = Presentation(str(path))
        self.page_width = cm(self.prs.slide_width)
        self.page_height = cm(self.prs.slide_height)

    def slide(self, page: int):
        return self.prs.slides[page - 1]

    def shapes(self, page: int):
        return list(self.slide(page).shapes)

    def all_shapes(self):
        for page_idx, slide in enumerate(self.prs.slides, start=1):
            for shape in slide.shapes:
                yield page_idx, shape

    def text_shapes(self, page: int) -> list:
        return [shape for shape in self.shapes(page) if has_text(shape)]

    def text(self, page: int) -> str:
        return "\n".join(shape.text for shape in self.text_shapes(page))

    def find_text(self, page: int, needle: str, fuzzy_space: bool = True) -> list:
        if fuzzy_space:
            target = norm_text(needle)
            return [s for s in self.text_shapes(page) if target in norm_text(s.text)]
        return [s for s in self.text_shapes(page) if needle in s.text]

    def find_text_regex(self, page: int, pattern: str) -> list:
        rx = re.compile(pattern)
        return [s for s in self.text_shapes(page) if rx.search(norm_text(s.text))]


def cm(value) -> float:
    return float(value) / EMU_PER_CM


def pt(value) -> Optional[float]:
    if value is None:
        return None
    return float(value) / PT_PER_EMU


def box(shape) -> Box:
    return Box(cm(shape.left), cm(shape.top), cm(shape.width), cm(shape.height))


def norm_text(text: str) -> str:
    return re.sub(r"\s+", "", (text or "").replace("", ""))


def has_text(shape) -> bool:
    return bool(getattr(shape, "has_text_frame", False) and norm_text(shape.text))


def in_range(value: float, low: float, high: float, tolerance: float = 0.0) -> bool:
    return low - tolerance <= value <= high + tolerance


def box_in(shape, left: float, top: float, right: float, bottom: float, tolerance: float = 0.0) -> bool:
    b = box(shape)
    return (
        in_range(b.left, left, right, tolerance)
        and in_range(b.top, top, bottom, tolerance)
        and b.right <= right + tolerance
        and b.bottom <= bottom + tolerance
    )


def center_in(shape, left: float, top: float, right: float, bottom: float, tolerance: float = 0.0) -> bool:
    b = box(shape)
    return in_range(b.cx, left, right, tolerance) and in_range(b.cy, top, bottom, tolerance)


def overlaps(a: Box, b: Box, margin: float = 0.0) -> bool:
    return not (
        a.right <= b.left + margin
        or b.right <= a.left + margin
        or a.bottom <= b.top + margin
        or b.bottom <= a.top + margin
    )


def intersection_area(a: Box, b: Box) -> float:
    width = max(0.0, min(a.right, b.right) - max(a.left, b.left))
    height = max(0.0, min(a.bottom, b.bottom) - max(a.top, b.top))
    return width * height


def within_page(inspector: PPTInspector, shape, tolerance: float = 0.05) -> bool:
    b = box(shape)
    return (
        b.left >= -tolerance
        and b.top >= -tolerance
        and b.right <= inspector.page_width + tolerance
        and b.bottom <= inspector.page_height + tolerance
    )


def shape_type_name(shape) -> str:
    try:
        return shape.shape_type.name
    except Exception:
        return str(shape.shape_type)


def is_picture(shape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE


def is_group(shape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.GROUP


def is_line(shape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.LINE


def is_auto_shape(shape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE


def is_rounded_rectangle(shape) -> bool:
    """办公软件中的“圆角矩形”自选图形（prstGeom = roundRect）。"""
    try:
        return shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    except Exception:
        return False


def is_editable_non_picture(shape) -> bool:
    return not is_picture(shape)


def get_rgb(color_obj) -> Optional[tuple[int, int, int]]:
    try:
        rgb = color_obj.rgb
        if rgb is None:
            return None
        return tuple(int(str(rgb)[i : i + 2], 16) for i in (0, 2, 4))
    except Exception:
        return None


def any_fill_rgb(shape) -> Optional[tuple[int, int, int]]:
    """读取形状的填充色，支持纯色/渐变。
    办公软件里对渐变填充也能看到一个主色调，故取渐变的首个色标作为“看到的颜色”。
    """
    solid = fill_rgb(shape)
    if solid is not None:
        return solid
    # 渐变填充：读 spPr/gradFill 的首个 srgbClr
    try:
        sp = shape._element
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        for grad in sp.iter(f"{ns}gradFill"):
            for clr in grad.iter(f"{ns}srgbClr"):
                val = clr.get("val")
                if val and len(val) == 6:
                    return (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
    except Exception:
        pass
    return None


def fill_rgb(shape) -> Optional[tuple[int, int, int]]:
    try:
        return get_rgb(shape.fill.fore_color)
    except Exception:
        return None


def line_rgb(shape) -> Optional[tuple[int, int, int]]:
    try:
        return get_rgb(shape.line.fill.fore_color)
    except Exception:
        return None


def group_white_bg(shape) -> bool:
    """组合对象内含大面积白色背景（办公软件将组合背景存为组内的白色填充图形）。"""
    if not is_group(shape):
        return False
    gb = box(shape)
    for child in shape.shapes:
        cb = box(child)
        if is_white(fill_rgb(child)) and cb.width >= gb.width * 0.5 and cb.height >= gb.height * 0.5:
            return True
    return False


def first_run(shape):
    if not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                return run
    return None


def all_runs(shape):
    if not getattr(shape, "has_text_frame", False):
        return []
    runs = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text:
                runs.append(run)
    return runs


def font_name(shape) -> Optional[str]:
    run = first_run(shape)
    return run.font.name if run is not None else None


def font_size(shape) -> Optional[float]:
    run = first_run(shape)
    return pt(run.font.size) if run is not None else None


def font_rgb(shape) -> Optional[tuple[int, int, int]]:
    run = first_run(shape)
    return get_rgb(run.font.color) if run is not None else None


def font_bold(shape) -> Optional[bool]:
    run = first_run(shape)
    return run.font.bold if run is not None else None


def is_sans_font(name: Optional[str]) -> bool:
    if not name:
        return True
    lowered = name.lower()
    return any(item in lowered for item in SANS_FONTS)


def is_dark(rgb: Optional[tuple[int, int, int]]) -> bool:
    # 深色（黑色或深灰色）：整体偏暗即可，不限定必须是某个主题色
    if rgb is None:
        return True
    return sum(rgb) / 3 <= 130


def is_white(rgb: Optional[tuple[int, int, int]]) -> bool:
    # 白色：只要接近白（各通道都很亮）即可，不要求纯 255
    if rgb is None:
        return False
    return min(rgb) >= 215


def color_distance(rgb: Optional[tuple[int, int, int]], target: tuple[int, int, int]) -> float:
    if rgb is None:
        return 999.0
    return math.sqrt(sum((a - b) ** 2 for a, b in zip(rgb, target)))


# 以下颜色判定按“颜色类别”识别，只要办公软件里看上去属于该色系就算通过，
# 不再要求匹配某个具体主题色 RGB。


def is_light_blue(rgb):
    # 浅蓝：整体偏亮，且蓝通道不明显低于红/绿通道
    if rgb is None:
        return False
    r, g, b = rgb
    return min(r, g, b) >= 150 and b >= 170 and b + 15 >= r


def is_light_yellow(rgb):
    # 浅黄：红、绿通道均较亮，蓝通道明显更低
    if rgb is None:
        return False
    r, g, b = rgb
    return r >= 200 and g >= 190 and b + 20 <= min(r, g)


def is_light_pink(rgb):
    # 浅粉：红通道最高、整体偏亮
    if rgb is None:
        return False
    r, g, b = rgb
    return r >= 210 and g >= 160 and b >= 170 and r >= g


def is_light_green(rgb):
    # 浅绿：绿通道最高、整体偏亮
    if rgb is None:
        return False
    r, g, b = rgb
    return g >= 190 and g + 10 >= r and min(r, g, b) >= 140


def is_deep_blue(rgb):
    # 深蓝：蓝通道明显高于红通道，且整体偏暗
    if rgb is None:
        return False
    r, g, b = rgb
    return b >= 80 and b > r + 25 and b + 20 >= g and r <= 130


def is_orange_red(rgb):
    # 橙红：红通道最强，蓝通道偏低
    if rgb is None:
        return False
    r, g, b = rgb
    return r >= 150 and r > g and r > b + 30


def is_purple_red(rgb):
    # 紫红：红、蓝通道较高，绿通道明显偏低
    if rgb is None:
        return False
    r, g, b = rgb
    return r >= 100 and g <= 130 and b >= 40 and r > g and b > g - 30


def is_dark_green(rgb):
    # 墨绿：绿通道占主导，整体偏暗
    if rgb is None:
        return False
    r, g, b = rgb
    return g >= 60 and g >= r and g + 20 >= b and max(r, g, b) <= 190


def is_dark_gray_line(shape) -> bool:
    # 灰色线：R/G/B 三通道相差不大即算灰色，不再限定“深”灰或纯黑（排除纯白线）
    rgb = line_rgb(shape)
    if rgb is None:
        return False
    return max(rgb) - min(rgb) <= 50 and max(rgb) <= 230


DEFAULT_LINE_WIDTH_PT = 1.0


def line_width_pt(shape) -> float:
    # 未显式设置线宽时，按办公软件默认线宽 1.0 磅处理
    try:
        width = shape.line.width
        if width is None:
            return DEFAULT_LINE_WIDTH_PT
        return width.pt
    except Exception:
        return DEFAULT_LINE_WIDTH_PT


def line_ok(shape, min_pt: float = 1.5, max_pt: float = 2.5) -> bool:
    width = line_width_pt(shape)
    return is_line(shape) and is_dark_gray_line(shape) and min_pt <= width <= max_pt


def is_solid_line(shape) -> bool:
    """单实线：没有虚线样式（办公软件中实线的 dash_style 为 None 或 SOLID）。"""
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        return shape.line.dash_style in (None, MSO_LINE_DASH_STYLE.SOLID)
    except Exception:
        return True


def _line_arrow_ends(shape) -> tuple[bool, bool]:
    """返回 (线条首端有内置箭头, 尾端有内置箭头)。
    办公软件里“一体箭头”会把箭头样式直接挂在线条自身（spPr/ln 下的 headEnd/tailEnd）。"""
    try:
        sp = shape._element
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        ln = sp.find(f".//{ns}spPr/{ns}ln")
        if ln is None:
            return (False, False)

        def _has(tag):
            el = ln.find(f"{ns}{tag}")
            if el is None:
                return False
            t = el.get("type")
            return t is not None and t != "none"

        return (_has("headEnd"), _has("tailEnd"))
    except Exception:
        return (False, False)


def is_integrated_arrow(shape) -> bool:
    """一体箭头：办公软件里"直线/直线连接符/肘形连接符"等 LINE 型形状带内置端点箭头即为一体箭头。
    含盖用户所述"line"和"直线连接符"两类——两者在 python-pptx 中同属 MSO_SHAPE_TYPE.LINE。"""
    if not is_line(shape):
        return False
    head_arrow, tail_arrow = _line_arrow_ends(shape)
    return head_arrow or tail_arrow


def line_arrow_direction(shape) -> Optional[str]:
    """一体箭头朝向：'right' | 'left' | 'down' | 'up'。对无法判定方向的一体箭头返回 None。
    针对直线：由整体外框比例判定水平/竖直，由 xfrm 的 flipH/flipV 判定端点物理位置，
    由 headEnd/tailEnd 判定箭头在哪一端。"""
    if not is_integrated_arrow(shape):
        return None
    head_arrow, tail_arrow = _line_arrow_ends(shape)
    b = box(shape)
    try:
        sp = shape._element
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        xfrm = sp.find(f".//{ns}xfrm")
        flip_h = xfrm is not None and xfrm.get("flipH") == "1"
        flip_v = xfrm is not None and xfrm.get("flipV") == "1"
    except Exception:
        flip_h = flip_v = False
    # 水平线：宽度显著大于高度；默认（无 flipH）线段 head 在左端、tail 在右端
    if b.width >= max(b.height, 0.01) * 1.5:
        tail_on_right = not flip_h
        if tail_arrow:
            return 'right' if tail_on_right else 'left'
        if head_arrow:
            return 'left' if tail_on_right else 'right'
    # 竖直线：高度显著大于宽度；默认（无 flipV）线段 head 在上端、tail 在下端
    elif b.height >= max(b.width, 0.01) * 1.5:
        tail_on_bottom = not flip_v
        if tail_arrow:
            return 'down' if tail_on_bottom else 'up'
        if head_arrow:
            return 'up' if tail_on_bottom else 'down'
    return None


def is_down_arrowhead(shape) -> bool:
    """向下箭头：等腰三角形自选图形旋转到开口朝下（rot≈180°），或直线型一体箭头朝下。"""
    if is_integrated_arrow(shape) and line_arrow_direction(shape) == 'down':
        return True
    try:
        if shape.auto_shape_type != MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE:
            return False
        rot = shape.rotation % 360
        return abs(rot - 180) <= 20
    except Exception:
        return False


def is_left_arrowhead(shape) -> bool:
    """朝左箭头：等腰三角形自选图形旋转到尖端朝左（rot≈270°），或直线型一体箭头朝左。"""
    if is_integrated_arrow(shape) and line_arrow_direction(shape) == 'left':
        return True
    try:
        if shape.auto_shape_type != MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE:
            return False
        rot = shape.rotation % 360
        return abs(rot - 270) <= 20
    except Exception:
        return False


def is_right_arrowhead(shape) -> bool:
    """朝右箭头：等腰三角形自选图形旋转到尖端朝右（rot≈90°），或直线型一体箭头朝右。"""
    if is_integrated_arrow(shape) and line_arrow_direction(shape) == 'right':
        return True
    try:
        if shape.auto_shape_type != MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE:
            return False
        rot = shape.rotation % 360
        return abs(rot - 90) <= 20
    except Exception:
        return False


def is_up_arrowhead(shape) -> bool:
    """朝上箭头：等腰三角形自选图形尖端朝上（rot≈0°），或直线型一体箭头朝上。"""
    if is_integrated_arrow(shape) and line_arrow_direction(shape) == 'up':
        return True
    try:
        if shape.auto_shape_type != MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE:
            return False
        rot = shape.rotation % 360
        return rot <= 20 or rot >= 340
    except Exception:
        return False


def shape_text_matches(shape, *needles: str) -> bool:
    content = norm_text(shape.text if has_text(shape) else "")
    return all(norm_text(needle) in content for needle in needles)


def has_any_text(page_shapes: Iterable, needle: str) -> bool:
    target = norm_text(needle)
    return any(has_text(shape) and target in norm_text(shape.text) for shape in page_shapes)


def find_state_box(inspector: PPTInspector, page: int, label: str):
    matches = [shape for shape in inspector.text_shapes(page) if norm_text(label) in norm_text(shape.text)]
    return matches[0] if matches else None


def detail_box(shape) -> str:
    b = box(shape)
    return f"{shape_type_name(shape)} x={b.left:.2f}, y={b.top:.2f}, w={b.width:.2f}, h={b.height:.2f}cm"


def evaluate_dimension1(inspector: PPTInspector) -> list[RuleResult]:
    results: list[RuleResult] = []
    suffix_ok = inspector.path.suffix.lower() == ".pptx"
    opened = len(inspector.prs.slides) > 0
    results.append(RuleResult("维度1-文件格式与可打开", 0, suffix_ok and opened, f"扩展名={inspector.path.suffix}, 页数={len(inspector.prs.slides)}"))

    page_count_ok = len(inspector.prs.slides) == 33
    results.append(RuleResult("维度1-页数为33页", 0, page_count_ok, f"实际页数={len(inspector.prs.slides)}"))

    return results


def check_nav_and_logo(inspector: PPTInspector) -> tuple[bool, str]:
    shapes = inspector.shapes(19)
    # 细则点1：距上0.4cm–2.8cm范围内出现「背景为白色」的「导航栏组合对象」
    nav_groups = [s for s in shapes if is_group(s) and 0.4 <= box(s).top <= 2.8 and group_white_bg(s)]
    # 细则点2：左上角距左2cm–5.5cm、距上0.4cm–2.8cm范围内出现 Logo 图片或 Logo 图形（排除导航栏组合本身）
    logo = [s for s in shapes if s not in nav_groups and (is_picture(s) or is_auto_shape(s) or is_group(s))
            and 2.0 <= box(s).left <= 5.5 and 0.4 <= box(s).top <= 2.8]
    # 细则点3：导航栏和Logo不被流程图遮挡（流程图=状态框Q1-Q4、条件标注、连接线）
    flow_keys = ("Q1", "Q2", "Q3", "Q4", "X=1", "M=0", "T↑", "Y=1", "N=0", "R=1", "Z=1", "P=0")
    flow_shapes = [s for s in shapes if is_line(s) or (has_text(s) and any(k in norm_text(s.text) for k in flow_keys))]
    protected = [box(s) for s in nav_groups + logo]
    not_blocked = all(not overlaps(box(f), pb) for f in flow_shapes for pb in protected)
    ok = bool(nav_groups and logo and not_blocked)
    return ok, f"白底导航组合={len(nav_groups)}, Logo候选={len(logo)}, 未被流程图遮挡={not_blocked}"


def check_old_body_cleanup(inspector: PPTInspector) -> tuple[bool, str]:
    shapes = inspector.shapes(19)
    # 细则点1：距左15cm–27cm、距上11cm–18cm范围内没有出现旧正文文本
    old_text_region = Box(15, 11, 12, 7)  # left 15→27, top 11→18
    old_text = [s for s in inspector.find_text(19, "项目推进以‘可验证、可解释、可维护’为原则")
                if overlaps(box(s), old_text_region)]
    # 细则点2：距左0.8cm–11cm、距上3.5cm–11cm范围内没有出现「与流程图无关的大面积装饰图组」
    #   装饰图组 = 图片/图形/组合等非文本视觉元素；办公软件中以其外框中心落在该范围内判定“出现在范围内”
    flow_keys = ("Q1", "Q2", "Q3", "Q4", "X=1", "M=0", "T↑", "Y=1", "N=0", "R=1", "Z=1", "P=0", "设备指示状态切换流程图")
    old_decor = [s for s in shapes
                 if (is_picture(s) or is_group(s) or is_auto_shape(s))            # 装饰图组（图片/图形/组合）
                 and center_in(s, 0.8, 3.5, 11.0, 11.0)                           # 出现在该范围内
                 and box(s).area >= 8                                             # 大面积
                 and not is_line(s)
                 and not (has_text(s) and any(k in norm_text(s.text) for k in flow_keys))]  # 与流程图无关
    ok = not old_text and not old_decor
    return ok, f"旧正文={len(old_text)}, 与流程图无关的大面积装饰图组={len(old_decor)}"


def check_flow_overall(inspector: PPTInspector) -> tuple[bool, str]:
    shapes = inspector.shapes(19)
    state_labels = ["Q1", "Q2", "Q3", "Q4"]
    conds = ["X=1", "M=0", "T↑", "Y=1", "N=0", "R=1", "Z=1", "P=0"]

    # ---- 可单独选中编辑 & 非截图/非图片（细则点：4状态框+8条件标注+标题+箭头线条均可编辑，且不是图片）----
    state_boxes = [find_state_box(inspector, 19, label) for label in state_labels]
    cond_shapes = [inspector.find_text(19, c) for c in conds]
    title = inspector.find_text(19, "设备指示状态切换流程图")
    lines = [s for s in shapes if is_line(s)]
    editable = all(state_boxes) and all(cond_shapes) and bool(title) and len(lines) >= 8
    # 状态框/条件标注/标题/箭头线条本身不是图片（文字不是嵌入式图片、箭头不是图片、状态框不是截图）
    flow_editable_shapes = [s for s in state_boxes + [c[0] for c in cond_shapes if c] + (title or []) if s] + lines
    none_is_picture = not any(is_picture(s) for s in flow_editable_shapes)

    # ---- 流程图主体（4状态框+8条件标注+箭头线条，标题单独判定）----
    body = [s for s in shapes if (
        is_line(s)
        or (has_text(s) and any(k in norm_text(s.text) for k in state_labels + conds))
    )]
    if body:
        left = min(box(s).left for s in body)
        top = min(box(s).top for s in body)
        right = max(box(s).right for s in body)
        bottom = max(box(s).bottom for s in body)
        width = right - left
        height = bottom - top
        cx = (left + right) / 2
        # 细则点：主体位于距左8–25cm、距上2–19.5cm范围内
        region_ok = left >= 8.0 and right <= 25.0 and top >= 2.0 and bottom <= 19.5
        # 细则点：整体宽度14–16cm、高度16–17cm
        size_ok = 14.0 <= width <= 16.0 and 16.0 <= height <= 17.0
        # 细则点：水平中心位于页面中心线16.9cm附近
        center_ok = abs(cx - 16.9) <= 1.0
        # 细则点：上下左右留白均衡
        left_margin = left
        right_margin = inspector.page_width - right
        top_margin = top
        bottom_margin = inspector.page_height - bottom
        margin_ok = abs(left_margin - right_margin) <= 1.0 and abs(top_margin - bottom_margin) <= 1.0
    else:
        left = top = width = height = cx = 0.0
        region_ok = size_ok = center_ok = margin_ok = False

    # ---- 不遮挡顶部导航栏和Logo ----
    nav_groups = [s for s in shapes if is_group(s) and 0.4 <= box(s).top <= 2.8 and group_white_bg(s)]
    logo = [s for s in shapes if s not in nav_groups and (is_picture(s) or is_auto_shape(s) or is_group(s))
            and 2.0 <= box(s).left <= 5.5 and 0.4 <= box(s).top <= 2.8]
    protected = [box(s) for s in nav_groups + logo]
    flow_all = body + (title or [])
    not_blocked = all(not overlaps(box(f), pb) for f in flow_all for pb in protected)

    # ---- 标题（细则点：上方中间出现该文本，位置/字体/字号/加粗/颜色/居中）----
    title_shape = title[0] if title else None
    title_ok = False
    if title_shape:
        b = box(title_shape)
        size = font_size(title_shape)
        align_ok = any(p.alignment == PP_ALIGN.CENTER for p in title_shape.text_frame.paragraphs if norm_text(p.text))
        title_ok = (
            10.0 <= b.left and b.right <= 22.0          # 距左10cm–22cm
            and 2.8 <= b.top and b.bottom <= 4.1        # 距上2.8cm–4.1cm
            and is_sans_font(font_name(title_shape))    # 黑体/微软雅黑/相近无衬线
            and size is not None and 18 <= size <= 22   # 字号18–22磅
            and font_bold(title_shape) is True          # 加粗
            and is_dark(font_rgb(title_shape))          # 黑色或深灰色
            and align_ok                                # 水平居中
        )

    ok = editable and none_is_picture and region_ok and size_ok and center_ok and margin_ok and not_blocked and title_ok
    return ok, (f"可编辑={editable}, 非图片={none_is_picture}, 主体外接框=({left:.2f},{top:.2f},{width:.2f}x{height:.2f})cm, "
                f"区域={region_ok}, 尺寸={size_ok}, 中心cx={cx:.2f}({center_ok}), 留白均衡={margin_ok}, "
                f"未遮挡导航={not_blocked}, 标题={title_ok}")


def check_q_text_fonts(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则：流程图内的文本字体为黑体、微软雅黑或相近无衬线字体
    # 检查范围覆盖流程图全部文本对象：
    #   - 4 个状态框（Q1–Q4）
    #   - 8 个条件标注（X=1、M=0、T↑、Y=1、N=0、R=1、Z=1、P=0）
    #   - 标题"设备指示状态切换流程图"
    # 对每个对象遍历所有非空 run，任一 run 字体不符合即视为不合格。
    labels = ["Q1", "Q2", "Q3", "Q4"]
    conds = ["X=1", "M=0", "T↑", "Y=1", "N=0", "R=1", "Z=1", "P=0"]
    title_text = "设备指示状态切换流程图"

    detail = []
    bad = []

    # 收集流程图内每个待检查对象：(label, shape)
    targets: list[tuple[str, object]] = []

    for label in labels:
        shape = find_state_box(inspector, 19, label)
        if not shape:
            bad.append(f"{label}(未找到)")
            continue
        targets.append((label, shape))

    for cond in conds:
        found = inspector.find_text(19, cond)
        if not found:
            # 条件标注未找到不在此评分点扣分（其他评分点已覆盖存在性），仅跳过字体检查
            continue
        targets.append((cond, found[0]))

    title_shapes = inspector.find_text(19, title_text)
    if title_shapes:
        targets.append((title_text, title_shapes[0]))

    for label, shape in targets:
        runs = all_runs(shape)
        if not runs:
            continue
        # 记录该对象的首个 run 字体名用于详情输出
        first_name = runs[0].font.name
        detail.append(f"{label}={first_name}")
        # 逐 run 检查：任一 run 字体不符合无衬线要求即不合格
        for run in runs:
            name = run.font.name
            if not is_sans_font(name):
                snippet = run.text.strip().replace("\n", " ")[:10]
                bad.append(f"{label}[{snippet}]({name})")
                break

    return not bad, f"字体={detail}, 不合格={bad or '无'}"


def check_state_box(inspector: PPTInspector, label: str, bits_patterns: list[str], bg_text: str, line_text: str, region: tuple[float, float, float, float], fill_checker: Callable, line_checker: Callable) -> tuple[bool, str]:
    candidates = [s for s in inspector.text_shapes(19) if norm_text(label) in norm_text(s.text)]
    if not candidates:
        return False, f"未找到{label}文本框"
    shape = candidates[0]
    b = box(shape)
    # 细则点：圆角矩形
    is_round = is_rounded_rectangle(shape)
    # 细则点：位于给定位置范围内（距左X–Y、距上A–B 指圆角矩形左上角的位置范围）
    region_ok = region[0] <= b.left <= region[2] and region[1] <= b.top <= region[3]
    # 细则点：宽4cm–7cm，高2.5cm–3.5cm
    size_ok = 4 <= b.width <= 7 and 2.5 <= b.height <= 3.5
    # 细则点：填充为指定浅色、边线为指定深色（只看实际颜色）
    fill_ok = fill_checker(fill_rgb(shape))
    outline_ok = line_checker(line_rgb(shape))
    # 细则点：线宽2–3磅（未显式设置线宽时按默认1.0磅处理）
    width = line_width_pt(shape)
    line_width_ok = 2 <= width <= 3
    # 细则点：框内出现主状态文本、"主背景：X"、"线框：X"
    text_ok = any(pattern in norm_text(shape.text) for pattern in bits_patterns) and norm_text(bg_text) in norm_text(shape.text) and norm_text(line_text) in norm_text(shape.text)
    # 细则点：主状态文本字号17–21磅、加粗（读实际磅值/加粗）
    main_font_ok = any((run.font.size and 17 <= run.font.size.pt <= 21 and run.font.bold is True) for run in all_runs(shape) if label in run.text)
    # 细则点：说明文字（主背景/线框两行）字号11–13磅
    small_font_ok = any((run.font.size and 11 <= run.font.size.pt <= 13) for run in all_runs(shape) if "主背景" in run.text or "线框" in run.text)
    # 细则点：文本居中
    align_ok = all(paragraph.alignment == PP_ALIGN.CENTER for paragraph in shape.text_frame.paragraphs if norm_text(paragraph.text))
    ok = is_round and region_ok and size_ok and fill_ok and outline_ok and line_width_ok and text_ok and main_font_ok and small_font_ok and align_ok
    return ok, f"{detail_box(shape)}, 圆角矩形={is_round}, 区域={region_ok}, 尺寸={size_ok}, 填充={fill_rgb(shape)}, 边线={line_rgb(shape)}, 线宽={width}, 文本={text_ok}, 主字号17-21加粗={main_font_ok}, 说明字号11-13={small_font_ok}, 居中={align_ok}"


def line_shapes_in(inspector: PPTInspector, page: int, region: Box) -> list:
    return [s for s in inspector.shapes(page) if is_line(s) and overlaps(box(s), region)]


def has_integrated_arrow_in(shapes, region: Box) -> list:
    """区域内的“一体箭头”：LINE 型（直线或直线连接符）且自身带箭头端点（headEnd/tailEnd）。
    只要区域内出现这类一体箭头，就视为“箭头”这部分要求达标——不再限制方向/线宽/深灰色等。"""
    return [s for s in shapes if is_integrated_arrow(s) and overlaps(box(s), region)]


def check_q1_q2_arrow(inspector: PPTInspector) -> tuple[bool, str]:
    shapes = inspector.shapes(19)
    # 连接区域：Q1右侧向右、再向下到Q2上方
    region = Box(19.5, 5.2, 8.5, 4.2)

    lines = [s for s in shapes if is_line(s) and overlaps(box(s), region)]
    # 细则点：线条为深灰色或黑灰色单实线，线宽1.5–2.5磅（未显式设置线宽按默认1.0磅处理）
    def seg_ok(s):
        w = line_width_pt(s)
        return is_dark_gray_line(s) and is_solid_line(s) and 1.5 <= w <= 2.5
    # 细则点：水平段长度约2.5cm–4.5cm
    horizontal = [s for s in lines if 2.5 <= box(s).width <= 4.5 and box(s).height <= 0.2 and seg_ok(s)]
    # 细则点：竖直段长度约3cm–5cm
    vertical = [s for s in lines if 3.0 <= box(s).height <= 5.0 and box(s).width <= 0.2 and seg_ok(s)]
    # 细则点：包含一个90度折角（水平段与竖直段首尾相接，形成直角）
    corner = any(abs(box(h).right - box(v).left) <= 0.5 or abs(box(h).right - box(v).right) <= 0.5
                 for h in horizontal for v in vertical)
    # 细则点：末端为向下箭头（竖直段下端出现开口朝下的箭头）
    down_arrow = [s for s in shapes if is_down_arrowhead(s) and overlaps(box(s), region)]
    arrow_ok = any(box(a).top >= box(v).bottom - 1.0 for a in down_arrow for v in vertical) if vertical else False
    # 细则点：竖直段右侧出现文本“X=1”，字号13–15磅
    label = inspector.find_text(19, "X=1")
    label_ok = bool(label and 13 <= (font_size(label[0]) or 0) <= 15
                    and vertical and box(label[0]).left >= max(box(v).right for v in vertical) - 0.5)
    ok = bool(horizontal and vertical and corner and arrow_ok and label_ok)
    return ok, f"水平段={len(horizontal)}, 竖直段={len(vertical)}, 90度折角={corner}, 向下箭头={bool(down_arrow) and arrow_ok}, X=1={label_ok}"


def check_q2_q4_arrow(inspector: PPTInspector) -> tuple[bool, str]:
    shapes = inspector.shapes(19)
    # 细则点：从Q2左侧水平指向Q4右侧，位于距上9.4cm–10.8cm范围内
    # 细则点：深灰/黑灰单实线，线宽1.5–2.5磅
    def seg_ok(s):
        w = line_width_pt(s)
        return is_dark_gray_line(s) and is_solid_line(s) and 1.5 <= w <= 2.5
    lines = [s for s in shapes if is_line(s)
             and 9.4 <= box(s).top <= 10.8 and box(s).height <= 0.2   # 水平、位于该纵向范围
             and 5.0 <= box(s).width <= 7.0                            # 长度约5–7cm
             and seg_ok(s)]
    # 细则点：末端箭头朝左
    left_arrow = [s for s in shapes if is_left_arrowhead(s) and 9.4 <= box(s).cy <= 10.8]
    arrow_ok = any(box(a).left <= box(l).left + 1.0 for a in left_arrow for l in lines) if lines else False
    # 细则点：线条上方中间出现文本“M=0”，字号13–15磅
    label = inspector.find_text(19, "M=0")
    label_ok = False
    if label and lines:
        lb = box(label[0])
        line = lines[0]; lnb = box(line)
        label_ok = (13 <= (font_size(label[0]) or 0) <= 15
                    and lb.bottom <= lnb.top + 0.3               # 在线条上方
                    and lnb.left - 1.0 <= lb.cx <= lnb.right + 1.0)  # 中间
    ok = bool(lines and arrow_ok and label_ok)
    return ok, f"水平线={len(lines)}, 箭头朝左={bool(left_arrow) and arrow_ok}, M=0={label_ok}"


def check_q2_q3_arrow(inspector: PPTInspector) -> tuple[bool, str]:
    shapes = inspector.shapes(19)
    region = Box(19.5, 11.7, 8.5, 3.4)

    lines = [s for s in shapes if is_line(s) and overlaps(box(s), region)]
    # 细则点：深灰/黑灰色单实线，线宽1.5–2.5磅
    def seg_ok(s):
        w = line_width_pt(s)
        return is_dark_gray_line(s) and is_solid_line(s) and 1.5 <= w <= 2.5
    # 细则点：包含竖直段（长度约3cm–5cm）和水平段（长度约2.5cm–4.5cm）
    horizontal = [s for s in lines if 2.5 <= box(s).width <= 4.5 and box(s).height <= 0.2 and seg_ok(s)]
    vertical = [s for s in lines if 3.0 <= box(s).height <= 5.0 and box(s).width <= 0.2 and seg_ok(s)]
    # 细则点：水平段箭头朝左
    left_arrow = [s for s in shapes if is_left_arrowhead(s) and overlaps(box(s), region)]
    arrow_ok = any(box(a).left <= box(h).left + 1.0 for a in left_arrow for h in horizontal) if horizontal else False
    # 细则点：Q3右侧附近出现文本“T↑”，字号13–15磅
    label = inspector.find_text(19, "T↑")
    label_ok = bool(label and 13 <= (font_size(label[0]) or 0) <= 15)
    ok = bool(horizontal and vertical and arrow_ok and label_ok)
    return ok, f"水平段={len(horizontal)}, 竖直段={len(vertical)}, 箭头朝左={bool(left_arrow) and arrow_ok}, T↑={label_ok}"


def check_q4_q1_arrow(inspector: PPTInspector) -> tuple[bool, str]:
    shapes = inspector.shapes(19)
    region = Box(6.2, 5.2, 7.8, 4.2)

    lines = [s for s in shapes if is_line(s) and overlaps(box(s), region)]
    # 细则点：深灰/黑灰色单实线，线宽1.5–2.5磅
    def seg_ok(s):
        w = line_width_pt(s)
        return is_dark_gray_line(s) and is_solid_line(s) and 1.5 <= w <= 2.5
    # 细则点：水平段长度约2.5cm–4.5cm；竖直段长度约3cm–5cm
    horizontal = [s for s in lines if 2.5 <= box(s).width <= 4.5 and box(s).height <= 0.2 and seg_ok(s)]
    vertical = [s for s in lines if 3.0 <= box(s).height <= 5.0 and box(s).width <= 0.2 and seg_ok(s)]
    # 细则点：包含一个90度折角（竖直段与水平段端点相接）
    corner = any(abs(box(v).top - box(h).top) <= 0.5 or abs(box(v).bottom - box(h).top) <= 0.5
                 for v in vertical for h in horizontal)
    # 细则点：箭头朝右（水平段右端出现尖端朝右的箭头）
    right_arrow = [s for s in shapes if is_right_arrowhead(s) and overlaps(box(s), region)]
    arrow_ok = any(box(a).left >= box(h).right - 1.0 for a in right_arrow for h in horizontal) if horizontal else False
    # 细则点：上方水平段出现文本“Y=1”；左侧竖直段附近出现文本“N=0”；字号13–15磅
    y_label = inspector.find_text(19, "Y=1")
    n_label = inspector.find_text(19, "N=0")
    labels_ok = bool(y_label and n_label and 13 <= (font_size(y_label[0]) or 0) <= 15 and 13 <= (font_size(n_label[0]) or 0) <= 15)
    ok = bool(horizontal and vertical and corner and arrow_ok and labels_ok)
    return ok, f"水平段={len(horizontal)}, 竖直段={len(vertical)}, 90度折角={corner}, 箭头朝右={bool(right_arrow) and arrow_ok}, Y/N={labels_ok}"


def check_q4_q3_arrow(inspector: PPTInspector) -> tuple[bool, str]:
    shapes = inspector.shapes(19)
    region = Box(6.2, 11.7, 8.0, 3.4)

    lines = [s for s in shapes if is_line(s) and overlaps(box(s), region)]
    # 细则点：深灰/黑灰色单实线，线宽1.5–2.5磅
    def seg_ok(s):
        w = line_width_pt(s)
        return is_dark_gray_line(s) and is_solid_line(s) and 1.5 <= w <= 2.5
    # 细则点：包含水平段（长约2.5–4.5cm）和竖直段（长约3–5cm）
    horizontal = [s for s in lines if 2.5 <= box(s).width <= 4.5 and box(s).height <= 0.2 and seg_ok(s)]
    vertical = [s for s in lines if 3.0 <= box(s).height <= 5.0 and box(s).width <= 0.2 and seg_ok(s)]
    # 细则点：箭头朝右（水平段右端出现尖端朝右的箭头）
    right_arrow = [s for s in shapes if is_right_arrowhead(s) and overlaps(box(s), region)]
    arrow_ok = any(box(a).left >= box(h).right - 1.0 for a in right_arrow for h in horizontal) if horizontal else False
    # 细则点：Q4与Q3之间出现文本“R=1”，字号13–15磅
    label = inspector.find_text(19, "R=1")
    label_ok = bool(label and 13 <= (font_size(label[0]) or 0) <= 15)
    ok = bool(horizontal and vertical and arrow_ok and label_ok)
    return ok, f"水平段={len(horizontal)}, 竖直段={len(vertical)}, 箭头朝右={bool(right_arrow) and arrow_ok}, R=1={label_ok}"


def check_q4_loop(inspector: PPTInspector) -> tuple[bool, str]:
    shapes = inspector.shapes(19)
    region = Box(1.0, 10.2, 3.0, 2.4)

    lines = [s for s in shapes if is_line(s) and overlaps(box(s), region)]
    # 细则点：深灰/黑灰色单实线，线宽1.5–2.5磅
    def seg_ok(s):
        w = line_width_pt(s)
        return is_dark_gray_line(s) and is_solid_line(s) and 1.5 <= w <= 2.5
    # 细则点：向左伸出再折回的U形/矩形折线路径——含水平段（约1.4–1.8cm）与竖直段（约1.8–2.2cm）
    horizontal = [s for s in lines if 1.4 <= box(s).width <= 1.8 and box(s).height <= 0.2 and seg_ok(s)]
    vertical = [s for s in lines if 1.8 <= box(s).height <= 2.2 and box(s).width <= 0.2 and seg_ok(s)]
    # 细则点：箭头指向Q4左边缘（尖端朝右的箭头，位于Q4左侧附近）
    q4 = find_state_box(inspector, 19, "Q4")
    q4_left = box(q4).left if q4 else 3.81
    right_arrow = [s for s in shapes if is_right_arrowhead(s) and overlaps(box(s), region)
                   and box(s).right <= q4_left + 0.5]
    # 细则点：左侧路径旁出现文本“Z=1”，字号13–15磅
    label = inspector.find_text(19, "Z=1")
    label_ok = bool(label and 13 <= (font_size(label[0]) or 0) <= 15)
    ok = len(horizontal) >= 2 and bool(vertical) and bool(right_arrow) and label_ok
    return ok, f"水平段={len(horizontal)}, 竖直段={len(vertical)}, 箭头指向Q4={bool(right_arrow)}, Z=1={label_ok}"


def check_q3_loop(inspector: PPTInspector) -> tuple[bool, str]:
    shapes = inspector.shapes(19)
    region = Box(14.5, 15.8, 5.0, 2.0)

    lines = [s for s in shapes if is_line(s) and overlaps(box(s), region)]
    # 细则点：深灰/黑灰色单实线，线宽1.5–2.5磅
    def seg_ok(s):
        w = line_width_pt(s)
        return is_dark_gray_line(s) and is_solid_line(s) and 1.5 <= w <= 2.5
    # 细则点：向下再折回的U形折线——含水平段（约2.5–3.1cm）与竖直段（约1.0–1.5cm）
    horizontal = [s for s in lines if 2.5 <= box(s).width <= 3.1 and box(s).height <= 0.2 and seg_ok(s)]
    vertical = [s for s in lines if 1.0 <= box(s).height <= 1.5 and box(s).width <= 0.2 and seg_ok(s)]
    # 细则点：箭头朝上指向Q3底边（尖端朝上的箭头，位于Q3底边附近）
    q3 = find_state_box(inspector, 19, "Q3")
    q3_bottom = box(q3).bottom if q3 else 16.26
    up_arrow = [s for s in shapes if is_up_arrowhead(s) and overlaps(box(s), region)
                and box(s).top <= q3_bottom + 0.5]
    # 细则点：底部中间出现文本“P=0”，字号13–15磅
    label = inspector.find_text(19, "P=0")
    label_ok = bool(label and 13 <= (font_size(label[0]) or 0) <= 15)
    ok = len(horizontal) >= 2 and bool(vertical) and bool(up_arrow) and label_ok
    return ok, f"水平段={len(horizontal)}, 竖直段={len(vertical)}, 箭头朝上指向Q3={bool(up_arrow)}, P=0={label_ok}"


def check_flow_layout_quality(inspector: PPTInspector) -> tuple[bool, str]:
    def _box_distance(a: Box, b: Box) -> float:
        dx = max(a.left - b.right, b.left - a.right, 0)
        dy = max(a.top - b.bottom, b.top - a.bottom, 0)
        return math.hypot(dx, dy)

    shapes = inspector.shapes(19)
    labels = ["Q1", "Q2", "Q3", "Q4"]
    conds = ["X=1", "M=0", "T↑", "Y=1", "N=0", "R=1", "Z=1", "P=0"]
    states = [find_state_box(inspector, 19, label) for label in labels]
    lines = [s for s in shapes if is_line(s)]

    # 细则点1：4个状态框之间无明显重叠
    state_overlap = []
    for i, a in enumerate(states):
        for j, b in enumerate(states[i + 1 :], start=i + 1):
            if a and b and overlaps(box(a), box(b), margin=-0.1):
                state_overlap.append((labels[i], labels[j]))

    # 细则点2：箭头不穿过状态框正文文字（箭头线条不进入状态框内部）
    state_inner_boxes = [Box(box(s).left + 0.2, box(s).top + 0.2, max(0, box(s).width - 0.4), max(0, box(s).height - 0.4)) for s in states if s]
    arrows_through_states = sum(1 for line in lines for state_box in state_inner_boxes if overlaps(box(line), state_box, margin=-0.02))

    # 细则点3：条件标注与箭头距离0.1cm–0.5cm
    bad_label_distances = []
    for cond in conds:
        found = inspector.find_text(19, cond)
        if not found or not lines:
            bad_label_distances.append((cond, None))
            continue
        distance = min(_box_distance(box(found[0]), box(line)) for line in lines)
        if not 0.1 <= distance <= 0.5:
            bad_label_distances.append((cond, round(distance, 2)))

    # 流程图内的全部文本（状态框/条件标注/标题）
    flow_texts = [s for s in inspector.text_shapes(19)
                  if any(k in norm_text(s.text) for k in labels + conds + ["设备指示状态切换流程图"])]

    # 细则点4：所有文本均在页面可视范围内（未被页面边界裁切）
    page_ok = all(within_page(inspector, s, tolerance=0.05) for s in flow_texts)

    # 细则点5：没有被导航栏、Logo或页面边界裁切
    #   导航栏=顶部白底组合，Logo=左上图片/图形（与其它评分项一致的识别方式）
    nav_groups = [s for s in shapes if is_group(s) and 0.4 <= box(s).top <= 2.8 and group_white_bg(s)]
    logo = [s for s in shapes if s not in nav_groups and (is_picture(s) or is_auto_shape(s) or is_group(s))
            and 2.0 <= box(s).left <= 5.5 and 0.4 <= box(s).top <= 2.8]
    protected = [box(s) for s in nav_groups + logo]
    # 文本被裁切=其中心落入导航栏/Logo区域（办公软件中中心被压住即视为遮挡裁切）
    clipped = [norm_text(s.text)[:8] for s in flow_texts
               if any(pb.left <= box(s).cx <= pb.right and pb.top <= box(s).cy <= pb.bottom for pb in protected)]
    not_clipped = not clipped

    ok = not state_overlap and arrows_through_states == 0 and not bad_label_distances and page_ok and not_clipped
    return ok, (f"状态框重叠={state_overlap or '无'}, 箭头穿状态框={arrows_through_states}, "
                f"标注距离异常={bad_label_distances or '无'}, 文本在页面内={page_ok}, "
                f"被导航/Logo裁切={clipped or '无'}")


def check_section_titles(inspector: PPTInspector) -> tuple[bool, str]:
    # 说明：细则写第23/29页，但“04.问题讨论”“05.研究结论”实际分别在第22/28页，
    #       按“在办公软件上实际看到的”为准，取实际所在页。
    section_targets = [(9, "02.研究思路"), (15, "03.项目进展"), (22, "04.问题讨论"), (28, "05.研究结论")]
    thanks = (33, "感谢聆听！")
    # 与第3页“01.研究背景”位置一致或接近
    ref = inspector.find_text(3, "01.研究背景")
    ref_box = box(ref[0]) if ref else None
    bad = []

    def base_ok(shape) -> tuple[bool, Box, Optional[float]]:
        b = box(shape)
        size = font_size(shape)
        # 细则点：字体为黑体/微软雅黑/相近无衬线字体
        font_ok = is_sans_font(font_name(shape))
        # 细则点：字号52-56磅（读实际磅值）
        size_ok = size is not None and 52 <= size <= 56
        # 细则点：位置约距左2.2-4cm、距上6.6-9cm（“约”取0.05cm容差）
        pos_ok = 2.2 - 0.05 <= b.left <= 4.0 + 0.05 and 6.6 - 0.05 <= b.top <= 9.0 + 0.05
        return (font_ok and size_ok and pos_ok), b, size

    for page, text in section_targets:
        found = inspector.find_text(page, text)
        if not found:
            bad.append(f"P{page}未找到{text}")
            continue
        shape = found[0]
        ok_base, b, size = base_ok(shape)
        # 细则点：与第3页“01.研究背景”位置一致或接近
        near_ref = ref_box is None or (abs(b.left - ref_box.left) <= 0.8 and abs(b.top - ref_box.top) <= 0.9)
        if not (ok_base and near_ref):
            bad.append(f"P{page}:{text} 字体={font_name(shape)} 字号={size} 位置=({b.left:.2f},{b.top:.2f}) 近01={near_ref}")

    # “感谢聆听！”
    page, text = thanks
    found = inspector.find_text(page, text)
    if not found:
        bad.append(f"P{page}未找到{text}")
    else:
        shape = found[0]
        ok_base, b, size = base_ok(shape)
        # 细则点：位于页面左侧近二分之一处居中（左半区 + 段落水平居中）
        left_half = b.left <= inspector.page_width / 2
        centered = any(p.alignment == PP_ALIGN.CENTER for p in shape.text_frame.paragraphs if norm_text(p.text))
        if not (ok_base and left_half and centered):
            bad.append(f"P{page}:{text} 字体={font_name(shape)} 字号={size} 位置=({b.left:.2f},{b.top:.2f}) 左半={left_half} 居中={centered}")

    return not bad, "; ".join(bad) if bad else "章节标题均符合字体/字号/位置要求"


def check_section_subtitles(inspector: PPTInspector) -> tuple[bool, str]:
    # 说明：细则写第23/29页，但“围绕应用边界…”“对研究价值…”实际在第22/28页，
    #       按“在办公软件上实际看到的”为准，取实际所在页。
    targets = [
        (9, "采用分阶段推进方式，将需求分析、数据治理、模型构建和验证反馈整合为一条可落地路径。"),
        (15, "阶段工作围绕数据整理、模型开发、流程验证和材料沉淀展开，已形成可汇报的中期成果。"),
        (22, "围绕应用边界、技术风险和后续优化路径展开讨论，重点识别影响落地效果的关键变量。"),
        (28, "对研究价值、阶段成果、技术难点和后续方向进行归纳，形成完整结论。"),
    ]
    # 第三页参照说明文字“围绕多…价值。”
    ref = inspector.find_text(3, "围绕多源设备数据治理、异常识别与预测维护需求，说明研究的现实基础与应用价值。")
    ref_shape = ref[0] if ref else None
    ref_box = box(ref_shape) if ref_shape else None
    ref_font = font_name(ref_shape) if ref_shape else None
    ref_size = font_size(ref_shape) if ref_shape else None
    bad = []
    for page, text in targets:
        found = inspector.find_text(page, text)
        if not found:
            bad.append(f"P{page}未找到说明文字")
            continue
        shape = found[0]
        b = box(shape)
        size = font_size(shape)
        # 细则点：字体为黑体/微软雅黑/相近无衬线字体
        font_ok = is_sans_font(font_name(shape))
        # 细则点：字号11-15磅（读实际磅值）
        size_ok = size is not None and 11 <= size <= 15
        # 细则点：与第三页“围绕多…价值。”位置一致或接近
        pos_ok = ref_box is None or (abs(b.left - ref_box.left) <= 0.8 and abs(b.top - ref_box.top) <= 0.8)
        # 细则点：与第三页字体一致或接近（同为无衬线，若参照可读则字体名一致）
        font_same = ref_font is None or is_sans_font(ref_font) == font_ok
        # 细则点：与第三页字号一致或接近（差值≤1磅）
        size_same = ref_size is None or size is None or abs(size - ref_size) <= 1
        ok = font_ok and size_ok and pos_ok and font_same and size_same
        if not ok:
            bad.append(f"P{page}: 字体={font_name(shape)} 字号={size} 位置=({b.left:.2f},{b.top:.2f}) 近第3页={pos_ok and size_same}")
    return not bad, "; ".join(bad) if bad else "分节说明文字均符合字体/字号/位置继承要求"


def check_toc_added_title(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则点：目录页出现“-05-”和“研究结论”两个可编辑文本框
    minus05 = [s for s in inspector.find_text(2, "-05-") if getattr(s, "has_text_frame", False)]
    title = [s for s in inspector.text_shapes(2) if norm_text(s.text) == norm_text("研究结论")]
    desc = inspector.find_text(2, "总结研究价值、阶段成果与后续优化方向。")
    if not (minus05 and title and desc):
        return False, f"-05-={bool(minus05)}, 研究结论={bool(title)}, 描述={bool(desc)}"
    s_num, s_title, s_desc = minus05[0], title[0], desc[0]

    # 参照：其他目录序号 -01-~-04- 与其他目录标题（研究背景/思路/进展/讨论）
    ref_nums = [inspector.find_text(2, k)[0] for k in ["-01-", "-02-", "-03-", "-04-"] if inspector.find_text(2, k)]
    ref_titles = [s for s in inspector.text_shapes(2)
                  if norm_text(s.text) in [norm_text(k) for k in ["研究背景", "研究思路", "项目进展", "问题讨论"]]]

    # 字体颜色为白色（办公软件中目录文字继承母版白色；无显式非白覆盖即视为白色）
    def effective_white(shape) -> bool:
        rgb = font_rgb(shape)
        return is_white(rgb) if rgb is not None else True

    def close_size(shape, refs) -> bool:
        size = font_size(shape)
        if size is None or not refs:
            return size is not None
        return all(font_size(r) is not None and abs(size - font_size(r)) <= 1 for r in refs)

    # 细则点：“-05-”为黑体/微软雅黑/相近无衬线字体28磅；字号、颜色、位置层级与-01-~-04-一致或接近
    num_font_ok = is_sans_font(font_name(s_num))
    num_size = font_size(s_num)
    num_size_exact_ok = num_size is not None and abs(num_size - 28) < 0.01        # 显式校验28磅
    num_size_ok = num_size_exact_ok and close_size(s_num, ref_nums)               # 兼顾与-01-~-04-一致或接近
    num_color_ok = effective_white(s_num) and all(effective_white(r) for r in ref_nums)  # 颜色一致（白）
    num_level_ok = all(abs(box(s_num).top - box(r).top) <= 0.3 for r in ref_nums) if ref_nums else True  # 位置层级（同一行高）
    num_ok = num_font_ok and num_size_ok and num_color_ok and num_level_ok

    # 细则点：“研究结论”为汉仪元隆黑90W/相近无衬线字体24磅；字号、颜色、对齐方式与其他目录标题一致或接近
    title_font_ok = is_sans_font(font_name(s_title))
    title_size = font_size(s_title)
    title_size_exact_ok = title_size is not None and abs(title_size - 24) < 0.01  # 显式校验24磅
    title_size_ok = title_size_exact_ok and close_size(s_title, ref_titles)       # 兼顾与其他目录标题一致或接近
    title_color_ok = effective_white(s_title) and all(effective_white(r) for r in ref_titles)  # 颜色一致（白）
    def para_align(s):
        return next((p.alignment for p in s.text_frame.paragraphs if norm_text(p.text)), None)
    title_align_ok = all(para_align(s_title) == para_align(r) for r in ref_titles) if ref_titles else True  # 对齐方式一致
    title_ok = title_font_ok and title_size_ok and title_color_ok and title_align_ok

    # 细则点：“总结研究价值、阶段成果与后续优化方向。”为黑体/微软雅黑/相近无衬线字体13磅
    desc_font_ok = is_sans_font(font_name(s_desc))
    desc_size = font_size(s_desc)
    desc_size_ok = desc_size is not None and abs(desc_size - 13) < 0.01           # 显式校验13磅

    # 细则点：字体颜色全部为白色
    white_ok = effective_white(s_num) and effective_white(s_title) and effective_white(s_desc)

    ok = num_ok and title_ok and desc_font_ok and desc_size_ok and white_ok
    return ok, (f"序号={font_name(s_num)}/{font_size(s_num)}(层级{num_level_ok},字号{num_size_ok}), "
                f"标题={font_name(s_title)}/{font_size(s_title)}(对齐{title_align_ok},字号{title_size_ok}), "
                f"描述={font_name(s_desc)}/{desc_size}(字号{desc_size_ok}), 全白={white_ok}")


def check_toc_added_position(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则：第2页“-05-研究结论”目录项位于目录页面「右侧或下方空白区域」，
    #       整体宽度约4cm–7cm、高度约6cm–9cm。
    #   校验标准 = 办公软件实际所见：
    #     · 目录项 = “-05-”“研究结论”“总结研究价值…”三个可编辑文本框（或其所在组合）；
    #       用它们的外接框（min-left/min-top/max-right/max-bottom）判定位置与宽高，
    #       rubric未要求深蓝色圆角矩形，避免误杀合格目录项。
    #     · 位于右侧空白区域 = 外接框左边缘落在页面右半侧；
    #       位于下方空白区域 = 外接框上边缘落在页面下半侧；两者任一成立即算“位于”。
    # 定位“-05-研究结论”目录项的三段文字（含在组合内也能命中，因为 text_shapes 会展平组合）
    texts = [s for s in inspector.text_shapes(2)
             if any(k in norm_text(s.text) for k in ["-05-", "研究结论", "总结研究价值"])]
    if not texts:
        return False, "未找到-05-研究结论目录项文字"

    # 目录项外接框：三个文本框（若某个缺失则按已有项取并集）
    boxes = [box(s) for s in texts]
    left = min(b.left for b in boxes)
    top = min(b.top for b in boxes)
    right = max(b.right for b in boxes)
    bottom = max(b.bottom for b in boxes)
    width, height = right - left, bottom - top

    # 位于目录页面右侧或下方空白区域
    on_right = left >= inspector.page_width / 2
    on_bottom = top >= inspector.page_height / 2
    position_ok = on_right or on_bottom
    # 整体宽度约4cm–7cm
    width_ok = 4 <= width <= 7
    # 高度约6cm–9cm
    height_ok = 6 <= height <= 9

    ok = position_ok and width_ok and height_ok
    return ok, (f"目录项外接框=({left:.2f},{top:.2f},{width:.2f}x{height:.2f})cm, "
                f"右侧或下方={position_ok}, 宽4-7={width_ok}, 高6-9={height_ok}")


def check_toc_added_no_overlap(inspector: PPTInspector) -> tuple[bool, str]:
    shapes = inspector.shapes(2)
    # 新增目录项“-05-研究结论”的三个文本框：序号/标题/描述
    s_num = next((s for s in inspector.text_shapes(2) if norm_text(s.text) == "-05-"), None)
    s_title = next((s for s in inspector.text_shapes(2) if norm_text(s.text) == norm_text("研究结论")), None)
    s_desc = next((s for s in inspector.text_shapes(2) if "总结研究价值" in norm_text(s.text)), None)
    if not (s_num and s_title and s_desc):
        return False, f"-05-={bool(s_num)}, 研究结论={bool(s_title)}, 描述={bool(s_desc)}"
    new_shapes = [s_num, s_title, s_desc]

    # 已有四个目录项的同类参照
    ref_nums = [inspector.find_text(2, k)[0] for k in ["-01-", "-02-", "-03-", "-04-"] if inspector.find_text(2, k)]
    ref_titles = [s for s in inspector.text_shapes(2)
                  if norm_text(s.text) in [norm_text(k) for k in ["研究背景", "研究思路", "项目进展", "问题讨论"]]]
    ref_descs = [s for s in inspector.text_shapes(2)
                 if any(k in norm_text(s.text) for k in ["围绕设备运行数据", "完成样本整理", "聚焦技术难点", "识别落地风险"])]

    # 细则点1：内容文本未超出文本框（办公软件中 auto_size 适应文字 或 开启自动换行 => 文本不溢出）
    def not_overflow(s) -> bool:
        try:
            tf = s.text_frame
            from pptx.enum.text import MSO_AUTO_SIZE
            return tf.word_wrap is True or tf.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        except Exception:
            return True
    text_inside = all(not_overflow(s) for s in new_shapes)

    # 细则点2：与已有四个目录项保持同类间距（用标题的左边距等差判断，-05-延续同一间距）
    spacing_ok = True
    if len(ref_titles) >= 2 and s_title is not None:
        lefts = sorted(box(t).left for t in ref_titles)
        gaps = [b - a for a, b in zip(lefts, lefts[1:])]
        avg_gap = sum(gaps) / len(gaps)
        spacing_ok = abs(box(s_title).left - lefts[-1] - avg_gap) <= 0.5

    # 细则点3：同样文本框大小（各文本框与其同类参照尺寸一致或接近）
    def same_size(s, refs) -> bool:
        if not refs:
            return True
        rw = sum(box(r).width for r in refs) / len(refs)
        rh = sum(box(r).height for r in refs) / len(refs)
        return abs(box(s).width - rw) <= 0.5 and abs(box(s).height - rh) <= 0.3
    size_ok = same_size(s_num, ref_nums) and same_size(s_title, ref_titles) and same_size(s_desc, ref_descs)

    # 细则点4：不遮挡“目录：”、已有目录文字、右上角装饰图、底部背景文字
    catalog = [s for s in inspector.text_shapes(2) if "目录" in norm_text(s.text)]
    existing_texts = ref_nums + ref_titles + ref_descs
    top_right_pic = [s for s in shapes if is_picture(s) and box(s).left >= inspector.page_width / 2 and box(s).top <= 5]
    bottom_texts = [s for s in inspector.text_shapes(2) if box(s).top >= inspector.page_height * 0.75 and s not in new_shapes]
    protected = catalog + existing_texts + top_right_pic + bottom_texts
    blocked = []
    for n in new_shapes:
        nb = box(n)
        for p in protected:
            inter = intersection_area(nb, box(p))
            if inter > min(nb.area, box(p).area) * 0.15 and nb.area > 0.1:
                name = norm_text(p.text)[:8] if has_text(p) else "装饰图/图片"
                blocked.append((norm_text(n.text)[:6], name or "装饰图/图片"))
    not_blocked = not blocked

    ok = text_inside and spacing_ok and size_ok and not_blocked
    return ok, f"文本未溢出={text_inside}, 同类间距={spacing_ok}, 同样文本框大小={size_ok}, 遮挡={blocked or '无'}"


def check_cover_reporter(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则点：出现可编辑文本“汇报人：陈前康邮箱:896700205@qq.com”
    found = inspector.find_text_regex(1, r"汇报人[:：]陈前康邮箱[:：]?896700205@qq\.com")
    if not found:
        return False, "未找到汇报人陈前康邮箱文本"
    shape = found[0]
    b = box(shape)
    # 细则点：位于距左18cm–26cm、距上12.4cm–13.7cm范围内（文本中心落入该范围）
    region_ok = 18 <= b.cx <= 26 and 12.4 <= b.cy <= 13.7
    # 细则点：字体为黑体/微软雅黑/相近无衬线字体
    font_ok = is_sans_font(font_name(shape))
    # 细则点：字号13磅
    size = font_size(shape)
    size_ok = size is not None and 12.5 <= size <= 13.5

    # 参照：封面其他人员信息（指导老师、Research presentation 等）
    refs = [s for s in inspector.text_shapes(1)
            if s is not shape and any(k in norm_text(s.text) for k in ["指导老师", "周沐辞", "Researchpresentation"])]
    # 细则点：字体、字号、颜色与封面其他人员信息一致或接近
    font_same = all(is_sans_font(font_name(r)) == font_ok for r in refs) if refs else True
    size_same = True
    if refs and size is not None:
        size_same = all(font_size(r) is None or abs(size - font_size(r)) <= 1 for r in refs)
    def eff_rgb(s):
        return font_rgb(s)  # None 表示继承（与同为继承的其他信息视为一致）
    color_same = True
    if refs:
        mine = eff_rgb(shape)
        for r in refs:
            rc = eff_rgb(r)
            if mine is None or rc is None:
                continue
            if color_distance(mine, rc) > 60:
                color_same = False
    consistent = font_same and size_same and color_same

    # 细则点：添加字体未与其余字体产生重叠
    others = [s for s in inspector.text_shapes(1) if s._element is not shape._element]
    overlapped = []
    for o in others:
        inter = intersection_area(b, box(o))
        if inter > 0.05:
            overlapped.append(norm_text(o.text)[:14])
    no_overlap = not overlapped

    ok = bool(found) and region_ok and font_ok and size_ok and consistent and no_overlap
    return ok, (f"{detail_box(shape)}, 区域={region_ok}, 字体={font_name(shape)}, 字号={size}(13磅={size_ok}), "
                f"与他人一致={consistent}, 未重叠={no_overlap}{'' if no_overlap else '('+str(overlapped)+')'}")


def check_slide19_logo_missing(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则点：第19页距左2cm–5.5cm、距上0.4cm–2.6cm范围内没有出现Logo图片或Logo图形
    #   Logo=图片或图形（图片/自选图形/组合），以其左上角落入该范围判定“出现在范围内”
    shapes = inspector.shapes(19)
    nav_groups = [s for s in shapes if is_group(s) and 0.4 <= box(s).top <= 2.8 and group_white_bg(s)]
    found = [s for s in shapes if s not in nav_groups and (is_picture(s) or is_auto_shape(s) or is_group(s))
             and 2.0 <= box(s).left <= 5.5 and 0.4 <= box(s).top <= 2.6]
    # 扣分命中 = 没有Logo
    return not found, f"Logo候选={len(found)}"


def is_blue_green(rgb: Optional[tuple[int, int, int]]) -> bool:
    """蓝绿色（青/teal 一类）：绿、蓝通道明显，且红通道低于绿和蓝。"""
    if rgb is None:
        return False
    r, g, b = rgb
    return g >= 90 and b >= 90 and r + 25 <= g and r + 25 <= b


def check_slide19_nav_bad(inspector: PPTInspector) -> tuple[bool, str]:
    shapes = inspector.shapes(19)
    # 细则点A：距上0.4cm–2.8cm范围内出现「蓝绿色」导航栏组合对象
    def group_blue_green_bg(g) -> bool:
        if not is_group(g):
            return False
        gb = box(g)
        for c in g.shapes:
            cb = box(c)
            if is_blue_green(fill_rgb(c)) and cb.width >= gb.width * 0.5 and cb.height >= gb.height * 0.5:
                return True
        return False
    nav_groups = [s for s in shapes if is_group(s) and 0.4 <= box(s).top <= 2.8 and group_blue_green_bg(s)]

    # 细则点B：导航栏区域（距上0.4cm–2.8cm）出现流程图状态框、箭头、标题文字
    nav_region = Box(0.0, 0.4, inspector.page_width, 2.4)  # top 0.4 → 2.8
    flow_keys = ("Q1", "Q2", "Q3", "Q4", "设备指示状态切换流程图")
    flow_in_nav = [s for s in shapes if overlaps(box(s), nav_region) and (
        is_line(s)                                                         # 箭头
        or (has_text(s) and any(k in norm_text(s.text) for k in flow_keys))  # 状态框/标题文字
    )]

    # 扣分命中 = 没有蓝绿色导航栏组合，或导航区被流程图元素占用
    hit = (not nav_groups) or bool(flow_in_nav)
    return hit, f"蓝绿色导航组合={len(nav_groups)}, 导航区流程图元素={len(flow_in_nav)}"


def check_old_decoration_penalty(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则点：第19页距左0.8cm–11cm、距上3.5cm–11cm范围内出现「与流程图无关的大面积旧版装饰图形或图片对象」
    #   装饰=图片/图形/组合；出现在范围内=外框中心落入该范围；大面积；与流程图无关=不含Q/条件/标题文字
    flow_keys = ("Q1", "Q2", "Q3", "Q4", "X=1", "M=0", "T↑", "Y=1", "N=0", "R=1", "Z=1", "P=0", "设备指示状态切换流程图")
    found = [s for s in inspector.shapes(19)
             if (is_picture(s) or is_group(s) or is_auto_shape(s))          # 装饰图形或图片
             and center_in(s, 0.8, 3.5, 11.0, 11.0)                         # 出现在该范围内
             and box(s).area >= 8                                           # 大面积
             and not is_line(s)
             and not (has_text(s) and any(k in norm_text(s.text) for k in flow_keys))]  # 与流程图无关
    # 扣分命中 = 存在这样的装饰对象
    return bool(found), f"与流程图无关的大面积旧装饰={len(found)}"


def check_flow_big_picture_penalty(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则点1：出现宽度超过20cm、高度超过12cm的流程图图片对象
    big_pics = [s for s in inspector.shapes(19) if is_picture(s) and box(s).width > 20 and box(s).height > 12]
    # 细则点2：没有可单独编辑的“Q1 0 0”“Q2 01”“Q3 1 1”“Q4 10”文本框
    #   可单独编辑=独立的含文本框对象（办公软件中可点选编辑的文本框）
    def has_editable(*patterns) -> bool:
        return any(getattr(s, "has_text_frame", False) and any(p in norm_text(s.text) for p in patterns)
                   for s in inspector.text_shapes(19))
    editable_states = (has_editable("Q100") and has_editable("Q201")
                       and has_editable("Q311", "Q311") and has_editable("Q410"))
    # 扣分命中 = 有超大图片 且 缺少可编辑状态文本框
    hit = bool(big_pics) and not editable_states
    return hit, f"超大图片(>20x12cm)={len(big_pics)}, 可编辑Q状态文本框齐全={editable_states}"


def check_missing_text_penalty(inspector: PPTInspector, page: int, text: str) -> tuple[bool, str]:
    # 通用扣分：某页“没有出现”指定文本时命中。
    #   为贴合办公软件实际显示，匹配时去除所有空白，并对常见等价标点做归一：
    #   点号（. ． 。 ·）、短横（各种 dash）、冒号（: ：），只要显示出该文本即算出现。
    DASHES = "-‐‑‒–—―−－﹣"
    DOTS = ".．。·・"

    def norm_punct(t: str) -> str:
        t = re.sub(r"\s+", "", t or "")
        out = []
        for ch in t:
            if ch in DASHES:
                out.append("-")
            elif ch in DOTS:
                out.append(".")
            elif ch == "：":
                out.append(":")
            else:
                out.append(ch)
        return "".join(out)

    target = norm_punct(text)
    found = [s for s in inspector.text_shapes(page) if target in norm_punct(s.text)]
    # 扣分命中 = 没有出现
    return not found, f"找到={len(found)}"


def check_missing_any_condition(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则点：第19页“没有出现任一”条件标注文本（这8个标注一个都没出现时才扣分）
    conds = ["X=1", "M=0", "T↑", "Y=1", "N=0", "R=1", "Z=1", "P=0"]
    present = [c for c in conds if inspector.find_text(19, c)]
    # 扣分命中 = 一个都没出现
    hit = not present
    return hit, f"出现的标注={present or '无'}"


def check_toc_placeholder(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则点：第2页目录页出现该占位英文文本
    #   办公软件里该文本可能因换行/多个run/空格差异而不完全逐字符一致，
    #   故按“去除所有空白后包含”判定，标点保持不变。
    target = re.sub(r"\s+", "", "Click here to enter your text, change the color or size of the text.")
    found = [s for s in inspector.text_shapes(2) if target in re.sub(r"\s+", "", s.text or "")]
    # 扣分命中 = 出现该占位文本
    return bool(found), f"占位文本={len(found)}"


def check_toc_missing_05(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则点：第2页目录页“没有出现-05-文本”或“没有出现研究结论文本”（任一缺失即扣分）
    #   办公软件里“-”可能是多种短横字符（普通连字符/减号/en-dash/em-dash/全角连字符等），
    #   且可能有空格差异；只要办公软件显示为“-05-”即算出现——故对短横做等价归一后匹配。
    DASHES = "-‐‑‒–—―−－﹣"

    def norm_dash(t: str) -> str:
        t = re.sub(r"\s+", "", t or "")
        return "".join("-" if ch in DASHES else ch for ch in t)

    texts = [norm_dash(s.text) for s in inspector.text_shapes(2)]
    has_05 = any("-05-" in t for t in texts)
    has_conclusion = any(norm_text("研究结论") in re.sub(r"\s+", "", s.text or "") for s in inspector.text_shapes(2))

    missing = []
    if not has_05:
        missing.append("-05-")
    if not has_conclusion:
        missing.append("研究结论")
    # 扣分命中 = 任一缺失
    return bool(missing), f"缺失={missing or '无'}"


def check_original_toc_font_changed(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则：第2页“原字体变化”——序号“-01-”应为微软雅黑28、标题“研究背景”应为汉仪元隆黑90W(或相近无衬线)24、
    #       内容“围绕设备运行数据…”应为微软雅黑13；任一不满足即扣分。
    #   字体：按“办公软件实际看到的字体类别”判定（微软雅黑/黑体等无衬线皆可；标题额外接受汉仪元隆黑90W）。
    #   字号：按细则给的磅值判定（读实际磅值，±1 容差容纳四舍五入）。
    DASHES = "-‐‑‒–—―−－﹣"

    def norm_dash(t: str) -> str:
        t = re.sub(r"\s+", "", t or "")
        return "".join("-" if ch in DASHES else ch for ch in t)

    def find_by(inspector, needle, dash=False):
        for s in inspector.text_shapes(2):
            hay = norm_dash(s.text) if dash else re.sub(r"\s+", "", s.text or "")
            key = norm_dash(needle) if dash else re.sub(r"\s+", "", needle)
            if key in hay:
                return s
        return None

    def title_font_ok(name: Optional[str]) -> bool:
        # 汉仪元隆黑90W 或 相近无衬线字体
        if is_sans_font(name):
            return True
        if not name:
            return False
        low = name.lower()
        return "汉仪元隆黑" in name or "hy yuanlong" in low or "yuanlong" in low

    # (查找文本, 期望字号, 字体判定, 是否短横归一, 名称)
    requirements = [
        ("-01-", 28, is_sans_font, True, "序号-01-"),
        ("研究背景", 24, title_font_ok, False, "标题研究背景"),
        ("围绕设备运行数据、故障预警需求与应用场景展开", 13, is_sans_font, False, "内容"),
    ]
    bad = []
    for needle, expected_size, font_checker, dash, label in requirements:
        shape = find_by(inspector, needle, dash=dash)
        if shape is None:
            bad.append(f"{label}:缺失")
            continue
        name = font_name(shape)
        size = font_size(shape)
        font_ok = font_checker(name)
        size_ok = size is not None and abs(size - expected_size) <= 1
        if not (font_ok and size_ok):
            bad.append(f"{label}: {name}/{size}(应≈{expected_size})")
    # 扣分命中 = 任一不满足
    return bool(bad), "; ".join(bad) if bad else "原目录字体/字号均符合"


def check_cover_old_reporter(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则点：第1页距左18cm–26cm、距上12.4cm–13.7cm范围内出现文本“汇报人：林泽书”
    #   文本匹配去空白并容忍中英文冒号；位置以文本中心落入该范围判定。
    def norm_colon(t: str) -> str:
        return re.sub(r"\s+", "", (t or "")).replace("：", ":")
    target = norm_colon("汇报人：林泽书")
    found = [s for s in inspector.text_shapes(1)
             if target in norm_colon(s.text)
             and 18 <= box(s).cx <= 26 and 12.4 <= box(s).cy <= 13.7]
    # 扣分命中 = 出现该旧汇报人文本
    return bool(found), f"旧汇报人林泽书={len(found)}"


def check_cover_missing_email(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则点：第1页没有出现文本“陈前康邮箱:896700205@qq.com”（未出现即扣分）
    #   匹配去空白并容忍中英文冒号，只要办公软件里显示出该内容即算出现。
    def norm_colon(t: str) -> str:
        return re.sub(r"\s+", "", (t or "")).replace("：", ":")
    target = norm_colon("陈前康邮箱:896700205@qq.com")
    found = [s for s in inspector.text_shapes(1) if target in norm_colon(s.text)]
    # 扣分命中 = 没有出现
    return not found, f"邮箱文本={len(found)}"


def check_cover_overlap_research(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则点：第1页文本“陈前康邮箱:896700205@qq.com”与前方“Research presentation”重叠（重叠即扣分）
    #   文本匹配去空白并容忍中英文冒号；重叠=两文本外框有实质相交。
    def norm_colon(t: str) -> str:
        return re.sub(r"\s+", "", (t or "")).replace("：", ":")
    email_key = norm_colon("陈前康邮箱:896700205@qq.com")
    research_key = norm_colon("Researchpresentation")
    email = [s for s in inspector.text_shapes(1) if email_key in norm_colon(s.text)]
    research = [s for s in inspector.text_shapes(1) if research_key in norm_colon(s.text)]
    overlapped = []
    for e in email:
        for r in research:
            inter = intersection_area(box(e), box(r))
            if inter > 0.05:   # 有实质交叠面积即视为重叠
                overlapped.append(round(inter, 3))
    # 扣分命中 = 存在重叠
    hit = bool(overlapped)
    return hit, f"重叠={hit}{'(交叠面积'+str(overlapped)+'cm²)' if overlapped else ''}"


def check_slide15_title_bad(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则：第15页标题“03.项目进展”字体为艺术体，或字号小于50磅或大于60磅（任一成立即扣分）。
    #   点号可能为半角/全角，做归一以贴合办公软件实际显示。
    DOTS = ".．。·・"
    def norm_dot(t: str) -> str:
        t = re.sub(r"\s+", "", t or "")
        return "".join("." if ch in DOTS else ch for ch in t)
    target = norm_dot("03.项目进展")
    found = [s for s in inspector.text_shapes(15) if target in norm_dot(s.text)]
    if not found:
        # 标题都不存在，谈不上“正常标题”，按异常计（与“没有标题”一致地命中扣分）
        return True, "未找到标题“03.项目进展”"
    shape = found[0]
    name = font_name(shape)
    size = font_size(shape)
    # 细则点：字体为艺术体（= 非黑体/微软雅黑/相近无衬线字体）
    art_font = not is_sans_font(name)
    # 细则点：字号小于50磅或大于60磅
    size_bad = size is None or size < 50 or size > 60
    # 扣分命中 = 任一成立
    hit = art_font or size_bad
    return hit, f"字体={name}(艺术体={art_font}), 字号={size}(超出50-60={size_bad})"


def check_slide3_missing_required(inspector: PPTInspector) -> tuple[bool, str]:
    # 细则：第3页没有出现“01.研究背景”“PART-”“围绕多源设备数据治理、异常识别与预测维护
    #       需求，说明研究的现实基础与应用价值。”这三组文本时扣分（三组中任一缺失即命中）。
    #   校验标准 = 办公软件上实际看到的：只要页面上显示出该文本即算“出现”，
    #   不限定其所在文本框、字体、字号、颜色等；仅去空白并对点号/短横/冒号做等价归一，
    #   以容忍办公软件中同形标点的不同编码与文字被拆分成多段/多个文本框的情况。
    DASHES = "-‐‑‒–—―−－﹣"
    DOTS = ".．。·・"

    def norm_punct(t: str) -> str:
        t = re.sub(r"\s+", "", t or "")
        out = []
        for ch in t:
            if ch in DASHES:
                out.append("-")
            elif ch in DOTS:
                out.append(".")
            elif ch == "：":
                out.append(":")
            else:
                out.append(ch)
        return "".join(out)

    # 收集第3页上办公软件可见的全部文字（含组合内嵌套文本），拼成整页文本，
    # 使“出现”判定与实际显示一致（不受文字被拆分到不同文本框/组合内的影响）。
    def collect_text(shapes, acc: list[str]) -> None:
        for shape in shapes:
            if is_group(shape):
                collect_text(shape.shapes, acc)
            elif getattr(shape, "has_text_frame", False) and shape.text:
                acc.append(shape.text)

    pieces: list[str] = []
    collect_text(inspector.shapes(3), pieces)
    page_text = norm_punct("".join(pieces))

    req = ["01.研究背景", "PART-", "围绕多源设备数据治理、异常识别与预测维护需求，说明研究的现实基础与应用价值。"]
    missing = [text for text in req if norm_punct(text) not in page_text]
    # 扣分命中 = 有任一缺失
    return bool(missing), f"缺失={missing or '无'}"


def evaluate_dimension2(inspector: PPTInspector) -> list[RuleResult]:
    rules = _dimension2_rules()
    results = []
    for name, score, checker in rules:
        try:
            hit, detail = checker(inspector)
        except Exception as exc:
            hit, detail = False, f"检测异常：{exc.__class__.__name__}: {exc}"
        results.append(RuleResult(name, score, hit, detail))
    return results


def validate_pptx_open(path: Path) -> Optional[str]:
    if not path.exists():
        return f"文件不存在：{path}"
    if path.suffix.lower() != ".pptx":
        return f"文件格式不是 .pptx：{path.suffix}"
    if not zipfile.is_zipfile(path):
        return "PPTX不是有效ZIP包，无法正常打开"
    try:
        Presentation(str(path))
    except Exception as exc:
        return f"python-pptx无法打开文件：{exc}"
    return None


SCRIPT_ID = "047"


def _locate_document(dir_path: Path) -> Optional[Path]:
    """在给定目录里定位待评估的 PPTX 文档。

    规则：目录下所有 .pptx 文件中挑第一个（跳过临时文件如 ~$xxx.pptx）。
    找不到时返回 None。
    """
    candidates = [
        p for p in sorted(dir_path.iterdir())
        if p.is_file()
        and p.suffix.lower() == ".pptx"
        and not p.name.startswith("~$")
    ]
    return candidates[0] if candidates else None


def _empty_result(dir_path: str, file_name: str, status: str, error: Optional[str],
                  max_score: int) -> dict:
    return {
        "id": SCRIPT_ID,
        "file_name": file_name,
        "status": status,
        "error": error,
        "dim1_pass": False,
        "dim1_reason": error or "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": max_score,
    }


# 维度二规则清单：与 evaluate_dimension2 内部保持一致，用于计算 max_score
# 与在维度一未通过时也能返回结构完整的 dim2_items（全部未命中）。
def _dimension2_rules() -> list[tuple[str, int, Callable[[PPTInspector], tuple[bool, str]]]]:
    return [
        ("第19页顶部导航栏和Logo", 1, check_nav_and_logo),
        ("第19页中部旧正文清理区域", 1, check_old_body_cleanup),
        ("第19页流程图整体位置和可编辑性", 5, check_flow_overall),
        ("第19页Q1-Q4文本字体", 1, check_q_text_fonts),
        ("第19页Q1状态框", 3, lambda i: check_state_box(i, "Q1", ["Q100", "Q100"], "主背景：浅蓝", "线框：深蓝", (13, 4.0, 19, 8), is_light_blue, is_deep_blue)),
        ("第19页Q2状态框", 3, lambda i: check_state_box(i, "Q2", ["Q201"], "主背景：浅黄", "线框：橙红", (18, 8, 25, 12), is_light_yellow, is_orange_red)),
        ("第19页Q3状态框", 3, lambda i: check_state_box(i, "Q3", ["Q311", "Q31 1"], "主背景：浅粉", "线框：紫红", (13, 13, 19, 18), is_light_pink, is_purple_red)),
        ("第19页Q4状态框", 3, lambda i: check_state_box(i, "Q4", ["Q410"], "主背景：浅绿", "线框：墨绿", (8, 8, 15, 12), is_light_green, is_dark_green)),
        ("第19页Q1到Q2连接箭头", 1, check_q1_q2_arrow),
        ("第19页Q2到Q4连接箭头", 1, check_q2_q4_arrow),
        ("第19页Q2到Q3连接箭头", 1, check_q2_q3_arrow),
        ("第19页Q4到Q1连接箭头", 1, check_q4_q1_arrow),
        ("第19页Q4到Q3连接箭头", 1, check_q4_q3_arrow),
        ("第19页Q4左侧自循环箭头", 1, check_q4_loop),
        ("第19页Q3底部自循环箭头", 1, check_q3_loop),
        ("第19页流程图排版质量", 1, check_flow_layout_quality),
        ("第9/15/23/29/33页章节标题", 3, check_section_titles),
        ("第9/15/23/29页章节说明文字", 1, check_section_subtitles),
        ("第2页目录页新增标题5", 3, check_toc_added_title),
        ("第2页新增目录项位置", 1, check_toc_added_position),
        ("第2页新增目录项不遮挡", 1, check_toc_added_no_overlap),
        ("第1页汇报人内容修改", 1, check_cover_reporter),
        ("扣分：第19页没有Logo", -5, check_slide19_logo_missing),
        ("扣分：第19页导航栏缺失或被流程图占用", -5, check_slide19_nav_bad),
        ("扣分：第19页左侧旧装饰残留", -1, check_old_decoration_penalty),
        ("扣分：第19页流程图为大图片且缺少可编辑文本", -5, check_flow_big_picture_penalty),
        ("扣分：第19页没有流程图标题", -1, lambda i: check_missing_text_penalty(i, 19, "设备指示状态切换流程图")),
        ("扣分：第19页缺少任一条件标注", -1, check_missing_any_condition),
        ("扣分：第2页目录占位英文未删除", -1, check_toc_placeholder),
        ("扣分：第2页缺少-05-或研究结论", -1, check_toc_missing_05),
        ("扣分：第2页原字体变化", -3, check_original_toc_font_changed),
        ("扣分：第1页仍出现旧汇报人林泽书", -1, check_cover_old_reporter),
        ("扣分：第1页缺少陈前康邮箱", -1, check_cover_missing_email),
        ("扣分：第1页陈前康邮箱与Research presentation重叠", -1, check_cover_overlap_research),
        ("扣分：第9页没有02.研究思路", -1, lambda i: check_missing_text_penalty(i, 9, "02.研究思路")),
        ("扣分：第15页03.项目进展标题异常", -1, check_slide15_title_bad),
        ("扣分：第3页缺少基础文本", -1, check_slide3_missing_required),
    ]


def _max_score_of_rules(rules) -> int:
    # 维度二满分：所有正分项之和（扣分项最好情况为 0，不计入满分）
    return sum(score for _name, score, _fn in rules if score > 0)


def _empty_dim2_items(rules) -> list[dict]:
    return [
        {"rule": name, "max_delta": score, "delta": 0, "hit": False, "detail": ""}
        for name, score, _fn in rules
    ]


def evaluate(dir_path: str) -> dict:
    """统一入口：传入脚本所在目录，返回结构化评分字典。

    dir_path 语义：脚本所在目录的路径。脚本负责在该目录内定位并打开被评估的
    PPT/PPTX 文档（当前目录下第一个 .ppt/.pptx 文件）。
    """
    rules = _dimension2_rules()
    max_score = _max_score_of_rules(rules)

    try:
        directory = Path(dir_path)
        if not directory.exists() or not directory.is_dir():
            return _empty_result(dir_path, "", "error",
                                 f"目录不存在或不是目录：{dir_path}", max_score)

        path = _locate_document(directory)
        if path is None:
            return _empty_result(dir_path, "", "error",
                                 f"目录中未找到 .ppt/.pptx 文档：{dir_path}", max_score)

        file_name = path.name
        open_error = validate_pptx_open(path)
        if open_error:
            return _empty_result(dir_path, file_name, "error", open_error, max_score)

        inspector = PPTInspector(path)

        # ---- 维度一 ----
        dim1_results = evaluate_dimension1(inspector)
        dim1_pass = all(r.hit for r in dim1_results)
        dim1_reason = "" if dim1_pass else "; ".join(
            f"{r.name}：{r.detail}" for r in dim1_results if not r.hit
        )

        if not dim1_pass:
            return {
                "id": SCRIPT_ID,
                "file_name": file_name,
                "status": "ok",
                "error": None,
                "dim1_pass": False,
                "dim1_reason": dim1_reason,
                "dim2_items": _empty_dim2_items(rules),
                "total_score": 0,
                "max_score": max_score,
            }

        # ---- 维度二 ----
        dim2_items: list[dict] = []
        total = 0
        for name, score, checker in rules:
            try:
                hit, _detail = checker(inspector)
            except Exception as exc:
                hit, _detail = False, f"检测异常：{exc.__class__.__name__}: {exc}"
            delta = score if hit else 0
            total += delta
            dim2_items.append({
                "rule": name,
                "max_delta": score,
                "delta": delta,
                "hit": bool(hit),
                "detail": "",
            })

        return {
            "id": SCRIPT_ID,
            "file_name": file_name,
            "status": "ok",
            "error": None,
            "dim1_pass": True,
            "dim1_reason": "",
            "dim2_items": dim2_items,
            "total_score": total,
            "max_score": max_score,
        }

    except Exception as exc:
        # 兜底：任何未预期异常都归为 status=error，不外抛
        return _empty_result(dir_path, "", "error",
                             f"{exc.__class__.__name__}: {exc}", max_score)


if __name__ == "__main__":
    # 本地调试入口：接收脚本所在目录路径（默认取本脚本所在目录），
    # 调用 evaluate 并以 JSON 输出结果，方便脚本作者自测。
    target_dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(target_dir), ensure_ascii=False, indent=2))
