#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
自动评估《工业园中水回用处理流程_可编辑版.pptx》。

评分逻辑：
1. 先检查“维度1：可用与可修改性”。任一项不满足，直接输出 0 分，
   且不再检查维度2。
2. 维度1通过后，逐项检查“维度2：完成度评分细则”。命中任一得分点/
   扣分点即累计该项分值（本题给出的细则均为正分，代码也支持负分项）。
3. 最终只打印命中的评分项和最终得分。

实现说明：
- 使用 python-pptx 直接读取 PPT 结构，不依赖人工判断。
- 对颜色、尺寸、位置采用少量容差，以兼容 PowerPoint/WPS/主题字体导致的细微差异。
- 对“组合对象”“可编辑图标”等 PPT 中较难直接判定的要求，采用可自动化的等价判断：
  * “完整流程图组合对象”优先识别真正的 group shape；没有 group 时，以全部可见元素的
    外接矩形作为“可编辑流程图组合”的候选。
  * “可编辑图标”识别为非图片的矢量形状/自由曲线/小型形状，且位于对应文本框左侧附近。
"""

from __future__ import annotations

import json
import math
import os
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable, Iterable, Optional

_PPTX_IMPORT_ERROR: Optional[ImportError] = None
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError as exc:  # pragma: no cover
    Presentation = None
    MSO_SHAPE_TYPE = None
    PP_ALIGN = None
    MSO_ANCHOR = None
    _PPTX_IMPORT_ERROR = exc

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_active_theme_map: dict[str, str] = {}

EMU_PER_CM = 360000
EMU_PER_PT = 12700

TARGET_FILE = "工业园中水回用处理流程_可编辑版.pptx"
EXPECTED_TITLE = "工业园中水回用处理流程与水质保障路径"
EXPECTED_STEPS = [
    "厂区综合废水收集",
    "粗格栅与调节井",
    "pH均衡池",
    "混凝沉淀池",
    "多介质过滤器",
    "活性炭吸附塔",
    "超滤膜组件",
    "消毒接触池",
    "回用清水池",
    "冷却补水系统",
]
STEP_TEXT_ALIASES: dict[int, set[str]] = {}
AUX_SPECS = [
    ("酸碱调节", 1, (2.5, 4.0), (4.3, 5.0), (1.0, 1.3), (0.3, 0.5), "烧瓶"),
    ("PAC/PAM投加", 1, (3.5, 5.0), (4.3, 5.0), (1.3, 1.4), (0.3, 0.4), "药剂"),
    ("核心净化单元", 1, (6.4, 8.4), (4.3, 5.0), (1.4, 1.6), (0.3, 0.4), "水滴"),
    ("次氯酸钠投加", 3, (8.3, 10.1), (4.3, 5.0), (1.4, 1.6), (0.3, 0.4), "盾牌"),
]
BOTTOM_MODULES = [
    ("全过程在线监测", "pH/浊度/余氯/电导率", "盾牌加号"),
    ("水质达标保障", "满足回用水质标准", "剪贴板"),
    ("设备定期维护", "保障系统稳定运行", "齿轮"),
    ("数据记录追溯", "管理可视化、可追溯", "文档"),
]


@dataclass
class ShapeInfo:
    shape: object
    index: str
    parent_group: Optional[object] = None
    # 幻灯片坐标系下的实际 EMU（考虑了外层 group 的 chOff/chExt 变换）。
    # 未设置时回退到 python-pptx 直读的原始值。
    _left_emu: Optional[int] = None
    _top_emu: Optional[int] = None
    _width_emu: Optional[int] = None
    _height_emu: Optional[int] = None

    @property
    def left(self) -> float:
        v = self._left_emu if self._left_emu is not None else getattr(self.shape, "left", 0)
        return emu_to_cm(v)

    @property
    def top(self) -> float:
        v = self._top_emu if self._top_emu is not None else getattr(self.shape, "top", 0)
        return emu_to_cm(v)

    @property
    def width(self) -> float:
        v = self._width_emu if self._width_emu is not None else getattr(self.shape, "width", 0)
        return emu_to_cm(v)

    @property
    def height(self) -> float:
        v = self._height_emu if self._height_emu is not None else getattr(self.shape, "height", 0)
        return emu_to_cm(v)

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height

    @property
    def text(self) -> str:
        return getattr(self.shape, "text", "") or ""


@dataclass
class Context:
    ppt_path: Path
    prs: object
    slide: object
    shapes: list[ShapeInfo]
    slide_w: float
    slide_h: float
    chart: tuple[float, float, float, float]  # left, top, width, height
    step_boxes: list[ShapeInfo] = field(default_factory=list)
    number_circles: list[ShapeInfo] = field(default_factory=list)
    aux_boxes: list[ShapeInfo] = field(default_factory=list)
    bottom_boxes: list[ShapeInfo] = field(default_factory=list)


@dataclass
class CheckResult:
    points: int
    name: str
    passed: bool
    detail: str


@dataclass
class GateResult:
    ok: bool
    detail: str
    prs: Optional[object] = None


def emu_to_cm(value) -> float:
    try:
        return float(value) / EMU_PER_CM
    except Exception:
        return 0.0


def emu_to_pt(value) -> Optional[float]:
    if value is None:
        return None
    try:
        return float(value) / EMU_PER_PT
    except Exception:
        return None


def norm_text(text: str) -> str:
    return re.sub(r"\s+", "", text or "")


def nearly(value: float, target: float, tol: float = 0.12) -> bool:
    return abs(value - target) <= tol


def in_range(value: float, low: float, high: float, tol: float = 0.05) -> bool:
    lo, hi = sorted((low, high))
    return lo - tol <= value <= hi + tol


def rect_close(si: ShapeInfo, x_rng=None, y_rng=None, w_rng=None, h_rng=None, origin=(0, 0), tol=0.05) -> bool:
    x0, y0 = origin
    x, y = si.left - x0, si.top - y0
    w, h = abs(si.width), abs(si.height)
    checks = []
    if x_rng is not None:
        checks.append(in_range(x, x_rng[0], x_rng[1], tol))
    if y_rng is not None:
        checks.append(in_range(y, y_rng[0], y_rng[1], tol))
    if w_rng is not None:
        checks.append(in_range(w, w_rng[0], w_rng[1], tol))
    if h_rng is not None:
        checks.append(in_range(h, h_rng[0], h_rng[1], tol))
    return all(checks)


def _load_theme_map(prs) -> dict[str, str]:
    """从第一张幻灯片对应的 slideLayout → slideMaster → theme 中读取 clrScheme。

    返回 {schemeName(小写): 6位大写HEX}，例如 {'accent1': '008C80', 'bg1': 'FFFFFF', 'tx1':'000000', ...}。
    另外把 clrMap 里的 bg1/tx1/... 映射也一并解析成对应的 accent/dk/lt。
    """
    mapping: dict[str, str] = {}
    try:
        slide = prs.slides[0]
        master = slide.slide_layout.slide_master
        theme_part = None
        # slide master 关联的 theme part
        try:
            for rel in master.part.rels.values():
                if "theme" in rel.reltype:
                    theme_part = rel.target_part
                    break
        except Exception:
            theme_part = None
        if theme_part is None:
            return mapping
        # theme part 是通用 Part（不是 pptx 强类型 part），需从 blob 解析 XML
        from lxml import etree as _et
        try:
            root = _et.fromstring(theme_part.blob)
        except Exception:
            root = getattr(theme_part, "_element", None)
        if root is None:
            return mapping
        # clrScheme
        clr_scheme = root.find(f".//{{{_A_NS}}}clrScheme")
        if clr_scheme is not None:
            for child in clr_scheme:
                tag = child.tag.split("}", 1)[-1]
                # child 下面可能是 srgbClr / sysClr
                srgb = child.find(f"{{{_A_NS}}}srgbClr")
                sysc = child.find(f"{{{_A_NS}}}sysClr")
                if srgb is not None and srgb.get("val"):
                    mapping[tag.lower()] = srgb.get("val").upper()
                elif sysc is not None:
                    last = sysc.get("lastClr")
                    if last:
                        mapping[tag.lower()] = last.upper()
        # slideMaster clrMap：把 bg1/tx1/bg2/tx2 映射到 lt1/dk1/lt2/dk2
        try:
            clr_map = master.element.find(f".//{{{_P_NS}}}clrMap")
            if clr_map is not None:
                for k, v in clr_map.attrib.items():
                    if v and v.lower() in mapping:
                        mapping[k.lower()] = mapping[v.lower()]
        except Exception:
            pass
    except Exception:
        pass
    return mapping


def _hex_apply_mods(hex_str: str, mods: list[tuple[str, int]]) -> str:
    """应用 lumMod/lumOff/shade/tint 等修饰。mods 是 [(name, val_thousandth), ...]。"""
    try:
        r = int(hex_str[0:2], 16); g = int(hex_str[2:4], 16); b = int(hex_str[4:6], 16)
    except Exception:
        return hex_str
    # 转 HSL 简化处理：这里做近似——lumMod 乘亮度，lumOff 加亮度
    import colorsys
    h, l, s = colorsys.rgb_to_hls(r / 255, g / 255, b / 255)
    for name, val in mods:
        v = val / 100000.0
        if name == "lumMod":
            l *= v
        elif name == "lumOff":
            l += v
        elif name == "shade":
            l *= v
        elif name == "tint":
            l = l + (1 - l) * (1 - v)
        elif name == "alpha":
            pass  # 忽略透明度
    l = max(0.0, min(1.0, l))
    nr, ng, nb = colorsys.hls_to_rgb(h, l, s)
    return "{:02X}{:02X}{:02X}".format(int(round(nr * 255)), int(round(ng * 255)), int(round(nb * 255)))


def _resolve_color_element(elem) -> Optional[str]:
    """从 a:srgbClr / a:schemeClr / a:sysClr 等颜色元素解析为 6 位 HEX。"""
    if elem is None:
        return None
    tag = elem.tag.split("}", 1)[-1]
    mods: list[tuple[str, int]] = []
    for child in elem:
        ctag = child.tag.split("}", 1)[-1]
        val = child.get("val")
        if val is not None:
            try:
                mods.append((ctag, int(val)))
            except Exception:
                pass
    base: Optional[str] = None
    if tag == "srgbClr":
        base = (elem.get("val") or "").upper()
    elif tag == "sysClr":
        base = (elem.get("lastClr") or "").upper() or None
    elif tag == "schemeClr":
        name = (elem.get("val") or "").lower()
        base = _active_theme_map.get(name)
    if not base or len(base) != 6:
        return None
    if mods:
        base = _hex_apply_mods(base, mods)
    return base


def safe_rgb(color) -> Optional[str]:
    """优先直接取 .rgb；失败时回退到解析底层 XML 中的 schemeClr / sysClr 等元素。"""
    try:
        rgb = color.rgb
        if rgb is not None:
            return str(rgb).upper()
    except Exception:
        pass
    # 从底层 XML 中解析
    elem = None
    for attr in ("_xFill", "_xClr", "_element"):
        cand = getattr(color, attr, None)
        if cand is not None:
            elem = cand
            break
    if elem is None:
        # 试试常见的 _color 属性
        try:
            elem = color._color._element  # type: ignore[attr-defined]
        except Exception:
            elem = None
    if elem is None:
        return None
    for tag in ("srgbClr", "schemeClr", "sysClr"):
        found = elem.find(f".//{{{_A_NS}}}{tag}")
        if found is not None:
            resolved = _resolve_color_element(found)
            if resolved:
                return resolved
    return None


def _style_fill_rgb(si: ShapeInfo) -> Optional[str]:
    """当形状未显式设置 solidFill 时，从 p:style/a:fillRef 取出主题色作为兜底。"""
    try:
        refs = si.shape.element.xpath(".//p:style/a:fillRef")
    except Exception:
        return None
    for ref in refs:
        for tag in ("srgbClr", "schemeClr", "sysClr"):
            found = ref.find(f".//{{{_A_NS}}}{tag}")
            if found is not None:
                resolved = _resolve_color_element(found)
                if resolved:
                    return resolved
    return None


def _style_line_rgb(si: ShapeInfo) -> Optional[str]:
    """当形状未显式设置 line 颜色时，从 p:style/a:lnRef 取出主题色作为兜底。"""
    try:
        refs = si.shape.element.xpath(".//p:style/a:lnRef")
    except Exception:
        return None
    for ref in refs:
        for tag in ("srgbClr", "schemeClr", "sysClr"):
            found = ref.find(f".//{{{_A_NS}}}{tag}")
            if found is not None:
                resolved = _resolve_color_element(found)
                if resolved:
                    return resolved
    return None


def fill_rgb(si: ShapeInfo) -> Optional[str]:
    try:
        return safe_rgb(si.shape.fill.fore_color)
    except Exception:
        return None


def line_rgb(si: ShapeInfo) -> Optional[str]:
    try:
        return safe_rgb(si.shape.line.color)
    except Exception:
        return None


def fill_rgb_with_style(si: ShapeInfo) -> Optional[str]:
    """含 p:style/a:fillRef 主题兜底的填充色。"""
    rgb = fill_rgb(si)
    if rgb:
        return rgb
    return _style_fill_rgb(si)


def line_rgb_with_style(si: ShapeInfo) -> Optional[str]:
    """含 p:style/a:lnRef 主题兜底的描边色。"""
    rgb = line_rgb(si)
    if rgb:
        return rgb
    return _style_line_rgb(si)


def _uses_line_style_ref(si: ShapeInfo) -> bool:
    """形状是否走了主题线样式 (a:lnRef)，用于兜底判定颜色随主题定义。"""
    try:
        return bool(si.shape.element.xpath(".//a:lnRef"))
    except Exception:
        return False


def line_width_pt(si: ShapeInfo) -> Optional[float]:
    try:
        return emu_to_pt(si.shape.line.width)
    except Exception:
        return None


def prst(si: ShapeInfo) -> Optional[str]:
    try:
        els = si.shape.element.xpath(".//a:prstGeom")
        if els:
            return els[0].get("prst")
    except Exception:
        pass
    return None


def is_line_like(si: ShapeInfo) -> bool:
    """判定为线/箭头形状：PowerPoint 的直线 (line) 或直线箭头 (straightConnector1)。"""
    p = prst(si)
    return p in {"line", "straightConnector1"}


def has_head_arrow(si: ShapeInfo) -> bool:
    """带头部箭头。PowerPoint 直线的箭头端可以在 headEnd 或 tailEnd 元素上表达。"""
    try:
        return bool(si.shape.element.xpath(".//a:headEnd[@type]"))
    except Exception:
        return False


def has_tail_arrow(si: ShapeInfo) -> bool:
    try:
        return bool(si.shape.element.xpath(".//a:tailEnd[@type]"))
    except Exception:
        return False


def has_any_arrow(si: ShapeInfo) -> bool:
    return has_head_arrow(si) or has_tail_arrow(si)


def dash_style(si: ShapeInfo) -> str:
    try:
        return str(si.shape.line.dash_style or "").upper()
    except Exception:
        return ""


def is_dashed(si: ShapeInfo) -> bool:
    d = dash_style(si)
    return "DASH" in d or "DOT" in d


def is_solid_line(si: ShapeInfo) -> bool:
    return not is_dashed(si)


def hex_to_rgb(hex_color: Optional[str]) -> Optional[tuple[int, int, int]]:
    if not hex_color or len(hex_color) != 6:
        return None
    try:
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
    except Exception:
        return None


def color_distance(a: Optional[str], b: str) -> float:
    ra = hex_to_rgb(a)
    rb = hex_to_rgb(b)
    if not ra or not rb:
        return 999.0
    return math.sqrt(sum((x - y) ** 2 for x, y in zip(ra, rb)))


def color_near(a: Optional[str], targets: Iterable[str], tol: float = 55) -> bool:
    return any(color_distance(a, t) <= tol for t in targets)


def is_white(c: Optional[str]) -> bool:
    return color_near(c, ["FFFFFF", "FEFEFE", "F8F8F8"], 25)


def is_black(c: Optional[str]) -> bool:
    return color_near(c, ["000000", "111111", "1F1F1F"], 45)


def is_gray(c: Optional[str]) -> bool:
    rgb = hex_to_rgb(c)
    return bool(rgb and max(rgb) - min(rgb) < 28 and 80 <= sum(rgb) / 3 <= 210)


def is_dark_teal(c: Optional[str]) -> bool:
    return color_near(c, ["008C80", "00796B", "006D68", "007C7A", "006B60"], 65)


def is_teal_or_green(c: Optional[str]) -> bool:
    return color_near(c, ["008C80", "78CBBE", "2C7FA4", "00A99D", "16A085", "3AAFA9", "2E8B57"], 85)


def is_dark_blue(c: Optional[str]) -> bool:
    return color_near(c, ["063A74", "003A70", "004A80", "0B3D91", "1F4E79", "1E386B"], 75)


def is_module_main_blue(c: Optional[str]) -> bool:
    """底部模块主文本使用的稍浅蓝（视觉仍属‘深蓝色’范畴）。"""
    return is_dark_blue(c) or color_near(c, ["4874CB", "2E5496", "3B65B0"], 75)


def is_light_green(c: Optional[str]) -> bool:
    rgb = hex_to_rgb(c)
    if not rgb:
        return False
    r, g, b = rgb
    # 极浅绿：整体偏亮，绿分量突出（≥ 红且 ≥ 蓝），允许微暖或微冷色调
    return r >= 210 and g >= 230 and b >= 200 and g >= r - 10 and g >= b - 10


def is_light_blue(c: Optional[str]) -> bool:
    rgb = hex_to_rgb(c)
    if not rgb:
        return False
    r, g, b = rgb
    return r >= 220 and g >= 235 and b >= 240


def is_light_blue_gray(c: Optional[str]) -> bool:
    rgb = hex_to_rgb(c)
    if not rgb:
        return False
    r, g, b = rgb
    return r >= 210 and g >= 220 and b >= 225 and max(rgb) - min(rgb) <= 45


def shape_type_name(si: ShapeInfo) -> str:
    try:
        return str(si.shape.shape_type)
    except Exception:
        return ""


def is_vector_icon_candidate(si: ShapeInfo) -> bool:
    """可编辑图标的自动化近似：不是图片，尺寸较小，有矢量形状。"""
    if norm_text(si.text):
        return False
    if abs(si.width) > 0.55 or abs(si.height) > 0.55:
        return False
    if si.shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return False
    p = prst(si)
    return bool(p and p not in {"rect", "roundRect", "line"}) or si.shape.shape_type in {
        MSO_SHAPE_TYPE.FREEFORM,
        MSO_SHAPE_TYPE.AUTO_SHAPE,
    }


def all_text_runs(si: ShapeInfo):
    try:
        for para in si.shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text:
                    yield para, run
    except Exception:
        return


def font_names(si: ShapeInfo) -> list[str]:
    return [getattr(run.font, "name", None) or "" for _, run in all_text_runs(si)]


def font_sizes(si: ShapeInfo) -> list[float]:
    sizes = []
    for _, run in all_text_runs(si):
        try:
            if run.font.size:
                sizes.append(float(run.font.size.pt))
        except Exception:
            pass
    return sizes


def font_colors(si: ShapeInfo) -> list[Optional[str]]:
    return [safe_rgb(run.font.color) for _, run in all_text_runs(si)]


def font_bolds(si: ShapeInfo) -> list[Optional[bool]]:
    return [run.font.bold for _, run in all_text_runs(si)]


def has_song_font(si: ShapeInfo) -> bool:
    names = font_names(si)
    if not names:
        return True  # 主题字体/继承字体无法读取时不强行判错
    return any(("宋" in n) or ("SimSun" in n) or ("Song" in n) for n in names)


def has_times_font(si: ShapeInfo) -> bool:
    names = font_names(si)
    if not names:
        return True
    return any("Times" in n for n in names)


def sizes_in(si: ShapeInfo, low: float, high: float, tol: float = 0.3) -> bool:
    sizes = font_sizes(si)
    return bool(sizes) and all(low - tol <= s <= high + tol for s in sizes)


def all_bold(si: ShapeInfo) -> bool:
    vals = font_bolds(si)
    return bool(vals) and all(v is True for v in vals)


def any_bold(si: ShapeInfo) -> bool:
    return any(v is True for v in font_bolds(si))


def text_colors_match(si: ShapeInfo, pred: Callable[[Optional[str]], bool]) -> bool:
    colors = font_colors(si)
    return bool(colors) and all(pred(c) for c in colors if c is not None)


def paragraph_centered(si: ShapeInfo) -> bool:
    try:
        return all(p.alignment in (PP_ALIGN.CENTER, None) for p in si.shape.text_frame.paragraphs)
    except Exception:
        return False


def vertical_middle(si: ShapeInfo) -> bool:
    try:
        return si.shape.text_frame.vertical_anchor in (MSO_ANCHOR.MIDDLE, None)
    except Exception:
        return True


def _group_transform(group_shape) -> Optional[tuple[int, int, int, int, int, int]]:
    """读取 group 的 (off_x, off_y, ext_cx, ext_cy, chOff_x, chOff_y, chExt_cx, chExt_cy)。

    返回把子坐标 (cx, cy, cw, ch) 变换到父坐标所需的 (tx, ty, sx_num, sx_den, sy_num, sy_den)：
        px = tx + (cx - chOff_x) * sx_num / sx_den
        py = ty + (cy - chOff_y) * sy_num / sy_den
        pw = cw * sx_num / sx_den
        ph = ch * sy_num / sy_den
    """
    try:
        xfrm = group_shape.element.xpath(".//p:grpSpPr/a:xfrm")
        if not xfrm:
            return None
        x = xfrm[0]
        off = x.find("{http://schemas.openxmlformats.org/drawingml/2006/main}off")
        ext = x.find("{http://schemas.openxmlformats.org/drawingml/2006/main}ext")
        ch_off = x.find("{http://schemas.openxmlformats.org/drawingml/2006/main}chOff")
        ch_ext = x.find("{http://schemas.openxmlformats.org/drawingml/2006/main}chExt")
        if off is None or ext is None or ch_off is None or ch_ext is None:
            return None
        off_x = int(off.get("x")); off_y = int(off.get("y"))
        ext_cx = int(ext.get("cx")); ext_cy = int(ext.get("cy"))
        cho_x = int(ch_off.get("x")); cho_y = int(ch_off.get("y"))
        che_cx = int(ch_ext.get("cx")); che_cy = int(ch_ext.get("cy"))
        return (off_x, off_y, ext_cx, ext_cy, cho_x, cho_y, che_cx, che_cy)
    except Exception:
        return None


def _apply_group_transform(child, transform):
    """把子形状的原始 (left, top, width, height) EMU 通过父 group 变换映射到父坐标系。"""
    off_x, off_y, ext_cx, ext_cy, cho_x, cho_y, che_cx, che_cy = transform
    cx = getattr(child, "left", 0) or 0
    cy = getattr(child, "top", 0) or 0
    cw = getattr(child, "width", 0) or 0
    ch = getattr(child, "height", 0) or 0
    sx = (ext_cx / che_cx) if che_cx else 1.0
    sy = (ext_cy / che_cy) if che_cy else 1.0
    px = off_x + (cx - cho_x) * sx
    py = off_y + (cy - cho_y) * sy
    pw = cw * sx
    ph = ch * sy
    return int(round(px)), int(round(py)), int(round(pw)), int(round(ph))


def flatten_shapes(shapes, prefix="", parent_transforms: Optional[list] = None) -> list[ShapeInfo]:
    result: list[ShapeInfo] = []
    parent_transforms = parent_transforms or []
    for idx, shape in enumerate(shapes):
        index = f"{prefix}{idx}"

        # 依次应用祖先 group 的变换，把子坐标一路映射到幻灯片坐标。
        if parent_transforms:
            left_emu = getattr(shape, "left", 0) or 0
            top_emu = getattr(shape, "top", 0) or 0
            width_emu = getattr(shape, "width", 0) or 0
            height_emu = getattr(shape, "height", 0) or 0

            class _Tmp:
                pass

            tmp = _Tmp()
            tmp.left, tmp.top, tmp.width, tmp.height = left_emu, top_emu, width_emu, height_emu
            for tr in parent_transforms:
                tmp.left, tmp.top, tmp.width, tmp.height = _apply_group_transform(tmp, tr)
            si = ShapeInfo(
                shape=shape,
                index=index,
                _left_emu=tmp.left,
                _top_emu=tmp.top,
                _width_emu=tmp.width,
                _height_emu=tmp.height,
            )
        else:
            si = ShapeInfo(shape=shape, index=index)
        result.append(si)

        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                tr = _group_transform(shape)
                child_transforms = parent_transforms + [tr] if tr else parent_transforms
                for child in flatten_shapes(shape.shapes, prefix=f"{index}.", parent_transforms=child_transforms):
                    child.parent_group = shape
                    result.append(child)
        except Exception:
            pass
    return result


def bbox_of(shapes: Iterable[ShapeInfo]) -> Optional[tuple[float, float, float, float]]:
    xs = []
    ys = []
    for si in shapes:
        w, h = si.width, si.height
        # line may have negative extents in PPT XML; normalize endpoints
        xs.extend([si.left, si.left + w])
        ys.extend([si.top, si.top + h])
    if not xs or not ys:
        return None
    left, right = min(xs), max(xs)
    top, bottom = min(ys), max(ys)
    return (left, top, right - left, bottom - top)


def local_pos(si: ShapeInfo, ctx: Context) -> tuple[float, float, float, float]:
    cl, ct, _, _ = ctx.chart
    return si.left - cl, si.top - ct, abs(si.width), abs(si.height)


def find_text_shape(ctx: Context, exact_text: str) -> Optional[ShapeInfo]:
    target = norm_text(exact_text)
    candidates = [si for si in ctx.shapes if norm_text(si.text) == target]
    if candidates:
        return sorted(candidates, key=lambda s: (s.top, s.left))[0]
    return None


def find_text_contains(ctx: Context, *pieces: str) -> Optional[ShapeInfo]:
    target_pieces = [norm_text(p).lower() for p in pieces]
    candidates = []
    for si in ctx.shapes:
        txt = norm_text(si.text).lower()
        if txt and all(p in txt for p in target_pieces):
            candidates.append(si)
    if candidates:
        return sorted(candidates, key=lambda s: (s.top, s.left))[0]
    return None


def find_module_pair(ctx: Context, main: str, sub: str) -> tuple[Optional[ShapeInfo], Optional[ShapeInfo]]:
    """底部保障模块常拆成主/副两条独立的 rect 文本框（大小写不敏感）。
    返回 (main_shape, sub_shape)；找不到则相应位置为 None。"""
    main_key = norm_text(main).lower()
    sub_key = norm_text(sub).lower()
    main_shape: Optional[ShapeInfo] = None
    sub_shape: Optional[ShapeInfo] = None
    # 先尝试合并在同一形状里
    combined = find_text_contains(ctx, main, sub)
    if combined:
        return combined, combined
    for si in ctx.shapes:
        txt = norm_text(si.text).lower()
        if not txt:
            continue
        if main_key in txt and (main_shape is None or si.top < main_shape.top):
            main_shape = si
        if sub_key in txt and (sub_shape is None or si.top < sub_shape.top):
            sub_shape = si
    return main_shape, sub_shape


def line_segments(ctx: Context, color_pred: Optional[Callable[[Optional[str]], bool]] = None) -> list[ShapeInfo]:
    segs = []
    for si in ctx.shapes:
        if prst(si) != "line":
            continue
        if color_pred and not color_pred(line_rgb(si)):
            continue
        segs.append(si)
    return segs


def is_horizontal(si: ShapeInfo, max_abs_h=0.08) -> bool:
    return abs(si.height) <= max_abs_h and abs(si.width) > 0.05


def is_vertical(si: ShapeInfo, max_abs_w=0.08) -> bool:
    return abs(si.width) <= max_abs_w and abs(si.height) > 0.05


def line_points(si: ShapeInfo) -> tuple[tuple[float, float], tuple[float, float]]:
    x1, y1 = si.left, si.top
    x2, y2 = si.left + si.width, si.top + si.height
    # 尊重 xfrm 上的 flipH / flipV：图形被翻转时端点顺序也随之调换
    try:
        xfrm = si.shape.element.xpath(".//a:xfrm")
        if xfrm:
            if xfrm[0].get("flipH") in ("1", "true"):
                x1, x2 = x2, x1
            if xfrm[0].get("flipV") in ("1", "true"):
                y1, y2 = y2, y1
    except Exception:
        pass
    return (x1, y1), (x2, y2)


def points_left_to_right(si: ShapeInfo) -> bool:
    (x1, _), (x2, _) = line_points(si)
    return x2 > x1


def points_down(si: ShapeInfo) -> bool:
    (_, y1), (_, y2) = line_points(si)
    return y2 > y1


def points_up(si: ShapeInfo) -> bool:
    (_, y1), (_, y2) = line_points(si)
    return y2 < y1


def line_width_ok(si: ShapeInfo, target=0.75, tol=0.20) -> bool:
    w = line_width_pt(si)
    return w is not None and abs(w - target) <= tol


def detect_flow_chart(prs, slide, shapes: list[ShapeInfo]) -> tuple[float, float, float, float]:
    # 1) 优先选择 group shape。
    groups = [si for si in shapes if si.shape.shape_type == MSO_SHAPE_TYPE.GROUP]
    if groups:
        return max(groups, key=lambda s: abs(s.width * s.height)).left, max(groups, key=lambda s: abs(s.width * s.height)).top, abs(max(groups, key=lambda s: abs(s.width * s.height)).width), abs(max(groups, key=lambda s: abs(s.width * s.height)).height)

    # 2) 没有 group 时，用可见形状外接矩形。
    visible = [si for si in shapes if abs(si.width) > 0.01 or abs(si.height) > 0.01]
    box = bbox_of(visible)
    if box:
        return box
    return (0.0, 0.0, emu_to_cm(prs.slide_width), emu_to_cm(prs.slide_height))


def dimension1_gate(path: Path) -> GateResult:
    details = []
    ok = True

    if path.suffix.lower() != ".pptx":
        ok = False
        details.append(f"文件扩展名为 {path.suffix or '(无扩展名)'}，不是 .pptx")
    else:
        details.append(f"文件扩展名为 {path.suffix}：通过")

    if not path.exists():
        return GateResult(False, "文件不存在，无法打开。")

    prs = None
    try:
        prs = Presentation(str(path))
        details.append("文件可由 python-pptx 正常打开：通过")
    except Exception as exc:
        ok = False
        details.append(f"文件无法正常打开：{exc}")
        return GateResult(False, "；".join(details), None)

    slide_count = len(prs.slides)
    if slide_count != 1:
        ok = False
        details.append(f"幻灯片数量为 {slide_count}，不是 1 页")
    else:
        details.append("只包含 1 页幻灯片：通过")

    return GateResult(ok, "；".join(details), prs)


def build_context(path: Path, prs) -> Context:
    slide = prs.slides[0]
    shapes = flatten_shapes(slide.shapes)
    chart = detect_flow_chart(prs, slide, shapes)
    ctx = Context(
        ppt_path=path,
        prs=prs,
        slide=slide,
        shapes=shapes,
        slide_w=emu_to_cm(prs.slide_width),
        slide_h=emu_to_cm(prs.slide_height),
        chart=chart,
    )
    ctx.step_boxes = detect_step_boxes(ctx)
    ctx.number_circles = detect_number_circles(ctx)
    ctx.aux_boxes = [find_text_shape(ctx, spec[0]) for spec in AUX_SPECS if find_text_shape(ctx, spec[0])]
    ctx.bottom_boxes = [find_text_contains(ctx, main, sub) for main, sub, _ in BOTTOM_MODULES if find_text_contains(ctx, main, sub)]
    return ctx


def expected_step_text_ok(index: int, text: str) -> bool:
    actual = norm_text(text).lower()
    allowed = {norm_text(EXPECTED_STEPS[index]).lower()}
    allowed.update(norm_text(t).lower() for t in STEP_TEXT_ALIASES.get(index, set()))
    return actual in allowed or any(t and t in actual for t in allowed)


def step_box_geometry_candidates(ctx: Context) -> list[ShapeInfo]:
    cl, ct, _, _ = ctx.chart
    candidates = [
        si
        for si in ctx.shapes
        if prst(si) == "roundRect"
        and 0.40 - 0.08 <= si.left - cl <= 12.30 + 0.08
        and 1.50 - 0.08 <= si.top - ct <= 4.50 + 0.08
        and 0.80 - 0.08 <= abs(si.width) <= 0.90 + 0.08
        and 2.00 - 0.08 <= abs(si.height) <= 3.00 + 0.08
    ]
    return sorted(candidates, key=lambda s: s.left)


def associated_text_shape(ctx: Context, box: ShapeInfo) -> Optional[ShapeInfo]:
    """把与 step box 关联的文本框找出来：优先返回 box 自身（若已有文字），
    否则找一个中心落在 box 水平范围内、且垂直方向与 box 有交集的独立文本框。
    编号标签（纯数字文本，例如两位数"10"被单独拆到 ellipse 外的文本框里）不应被当作正文。"""
    if norm_text(box.text):
        return box
    best: Optional[ShapeInfo] = None
    for si in ctx.shapes:
        if si.index == box.index:
            continue
        try:
            if si.shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                continue
        except Exception:
            continue
        t = norm_text(si.text)
        if not t:
            continue
        # 跳过纯数字（1-10）的编号标签
        if t.isdigit():
            continue
        cx = si.left + abs(si.width) / 2
        # 中心 x 落在 box 内（放宽一点点，允许字体边距）
        if not (box.left - 0.08 <= cx <= box.right + 0.08):
            continue
        # 垂直方向与 box 有交集
        if si.bottom < box.top - 0.05 or si.top > box.bottom + 0.05:
            continue
        best = si
        break
    return best


def detect_step_boxes(ctx: Context) -> list[ShapeInfo]:
    # 优先按几何形状识别：10 个空的圆角矩形，位置和尺寸符合区间。
    geometry_candidates = step_box_geometry_candidates(ctx)
    if len(geometry_candidates) == 10:
        return geometry_candidates

    # 退化路径：根据框内文本匹配（用于历史版本 PPT 把文字嵌在圆角矩形内）。
    matched: list[ShapeInfo] = []
    used: set[str] = set()
    # 允许第 7/8 步文本重复：逐个按 x 坐标匹配最靠左的未使用候选。
    for i, expected in enumerate(EXPECTED_STEPS):
        allowed = {norm_text(expected)}
        allowed.update(norm_text(t) for t in STEP_TEXT_ALIASES.get(i, set()))
        candidates = [
            si
            for si in ctx.shapes
            if norm_text(si.text) in allowed
            and prst(si) == "roundRect"
            and si.index not in used
        ]
        if not candidates:
            candidates = [
                si
                for si in ctx.shapes
                if any(t and t in norm_text(si.text) for t in allowed)
                and prst(si) == "roundRect"
                and si.index not in used
            ]
        if not candidates:
            continue
        cand = sorted(candidates, key=lambda s: s.left)[0]
        used.add(cand.index)
        matched.append(cand)
    matched = sorted(matched, key=lambda s: s.left)
    return matched


def circle_associated_text(ctx: Context, circle: ShapeInfo) -> ShapeInfo:
    """圆内直接含数字则返回自身；否则返回最近的、包含数字的文本框。"""
    if norm_text(circle.text).isdigit():
        return circle
    cx = circle.left + circle.width / 2
    cy = circle.top + circle.height / 2
    best: Optional[ShapeInfo] = None
    best_dist = 0.6
    for si in ctx.shapes:
        try:
            if si.shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                continue
        except Exception:
            continue
        tt = norm_text(si.text)
        if not tt.isdigit():
            continue
        sx = si.left + si.width / 2
        sy = si.top + si.height / 2
        d = ((sx - cx) ** 2 + (sy - cy) ** 2) ** 0.5
        if d < best_dist:
            best_dist = d
            best = si
    return best if best is not None else circle


def is_white_or_no_fill(si: ShapeInfo) -> bool:
    """填充为白色 或 无显式填充（PowerPoint 中形状默认背景/无填充视觉上等同白色）。"""
    c = fill_rgb(si)
    if c is None:
        return True
    return is_white(c)


def detect_number_circles(ctx: Context) -> list[ShapeInfo]:
    # 直接命中：圆内含 1-10 数字
    inline_circles = [
        si
        for si in ctx.shapes
        if prst(si) == "ellipse" and norm_text(si.text) in {str(i) for i in range(1, 11)}
    ]

    # 若某些圆的数字在独立文本框里，用几何位置补齐：10 个直径 0.25-0.30cm 的圆，与关联文本框数字排序
    small_circles = [
        si
        for si in ctx.shapes
        if prst(si) == "ellipse"
        and in_range(abs(si.width), 0.22, 0.32, 0.0)
        and in_range(abs(si.height), 0.22, 0.32, 0.0)
    ]

    def circle_number(c: ShapeInfo) -> Optional[str]:
        t = norm_text(c.text)
        if t.isdigit():
            return t
        # 找中心距离最近、内部含数字的文本框
        best_txt = None
        best_dist = 0.6  # 允许最大 0.6cm 中心距
        cx, cy = c.left + c.width / 2, c.top + c.height / 2
        for si in ctx.shapes:
            try:
                if si.shape.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
                    continue
            except Exception:
                continue
            tt = norm_text(si.text)
            if not tt.isdigit():
                continue
            sx, sy = si.left + si.width / 2, si.top + si.height / 2
            d = ((sx - cx) ** 2 + (sy - cy) ** 2) ** 0.5
            if d < best_dist:
                best_dist = d
                best_txt = tt
        return best_txt

    labeled: list[tuple[int, ShapeInfo]] = []
    for c in small_circles:
        num = circle_number(c)
        if num and num.isdigit() and 1 <= int(num) <= 10:
            labeled.append((int(num), c))

    # 若能凑齐 1-10 十个不重复的编号，就用几何+关联文本的结果
    seen_nums = {n for n, _ in labeled}
    if len(labeled) >= 10 and seen_nums == set(range(1, 11)):
        # 按编号升序（若同编号有多个则按 left 位置）
        return [c for _, c in sorted(labeled, key=lambda kv: (kv[0], kv[1].left))]

    return sorted(inline_circles, key=lambda s: (int(norm_text(s.text)) if norm_text(s.text).isdigit() else 99, s.left))


def flow_chart_group_candidate(ctx: Context) -> Optional[ShapeInfo]:
    groups = [si for si in ctx.shapes if si.shape.shape_type == MSO_SHAPE_TYPE.GROUP]
    if not groups:
        return None
    return max(groups, key=lambda s: abs(s.width * s.height))


def flow_chart_region_ok(ctx: Context) -> bool:
    group = flow_chart_group_candidate(ctx)
    if not group:
        return False
    return (
        in_range(group.left, 10.5, 11.5, 0.0)
        and in_range(group.top, 5.5, 6.5, 0.0)
        and nearly(abs(group.width), 12.44, 0.01)
        and nearly(abs(group.height), 7.00, 0.01)
    )


def check_complete_group(ctx: Context) -> CheckResult:
    group = flow_chart_group_candidate(ctx)
    if not group:
        return CheckResult(5, "第1页完整流程图组合对象", False, "未找到办公软件中的组合对象")
    passed = flow_chart_region_ok(ctx)
    detail = f"识别到组合对象 left={group.left:.2f}cm, top={group.top:.2f}cm, width={abs(group.width):.2f}cm, height={abs(group.height):.2f}cm"
    return CheckResult(5, "第1页完整流程图组合对象", passed, detail)


def has_vertical_flip(si: ShapeInfo) -> bool:
    try:
        xfrm = si.shape.element.xpath(".//a:xfrm")
        return bool(xfrm and xfrm[0].get("flipV") in {"1", "true"})
    except Exception:
        return False


def rotation_is_level(si: ShapeInfo) -> bool:
    try:
        rotation = float(getattr(si.shape, "rotation", 0) or 0) % 360
        return rotation <= 0.5 or rotation >= 359.5
    except Exception:
        return True


def title_bg_shape_ok(si: ShapeInfo) -> bool:
    # PowerPoint/WPS 的“同侧圆角矩形”预设表示顶部左右角为圆角、底边为直线的标题背景。
    return prst(si) == "round2SameRect" and rotation_is_level(si) and not has_vertical_flip(si)


def check_title_bg(ctx: Context) -> CheckResult:
    group = flow_chart_group_candidate(ctx)
    if not group:
        return CheckResult(5, "第1页顶部标题背景", False, "未找到完整流程图组合对象，无法定位其顶部")
    candidates = [
        si for si in ctx.shapes
        if nearly(si.left, group.left, 0.05)
        and nearly(si.top, group.top, 0.05)
        and nearly(abs(si.width), 12.44, 0.05)
        and nearly(abs(si.height), 0.80, 0.05)
        and is_dark_teal(fill_rgb(si))
        and title_bg_shape_ok(si)
    ]
    passed = bool(candidates)
    detail = "找到位于组合对象顶部、12.44×0.80cm、深青绿色、顶部左右圆角且下边缘水平直线的标题背景" if passed else "未找到同时满足组合对象顶部、12.44×0.80cm、深青绿色、顶部左右圆角、下边缘水平直线的标题背景"
    return CheckResult(5, "第1页顶部标题背景", passed, detail)


def shape_center(si: ShapeInfo) -> tuple[float, float]:
    return si.left + abs(si.width) / 2, si.top + abs(si.height) / 2


def text_centered_in_bg(si: ShapeInfo, bg: ShapeInfo, tol: float = 0.20) -> bool:
    """文本框中心点落在标题背景中央：水平方向和垂直方向均在背景中心容差范围内。"""
    tx, ty = shape_center(si)
    bx, by = shape_center(bg)
    return abs(tx - bx) <= tol and abs(ty - by) <= tol


def check_title_text(ctx: Context) -> CheckResult:
    si = find_text_shape(ctx, EXPECTED_TITLE)
    if not si:
        return CheckResult(5, "第1页顶部标题文本", False, "未找到文本【工业园中水回用处理流程与水质保障路径】")

    # 找标题背景（复用 check_title_bg 的逻辑定位背景形状）
    group = flow_chart_group_candidate(ctx)
    bg_candidates = [
        bg for bg in ctx.shapes
        if group is not None
        and nearly(bg.left, group.left, 0.01)
        and nearly(bg.top, group.top, 0.01)
        and nearly(abs(bg.width), 12.44, 0.05)
        and nearly(abs(bg.height), 0.80, 0.05)
        and is_dark_teal(fill_rgb(bg))
    ] if group else []
    bg = bg_candidates[0] if bg_candidates else None

    # 细则各点
    text_ok = bool(norm_text(si.text) == norm_text(EXPECTED_TITLE))
    centered_ok = text_centered_in_bg(si, bg) if bg else False
    font_ok = has_song_font(si)
    size_ok = sizes_in(si, 12, 14)
    bold_ok = all_bold(si)
    color_ok = text_colors_match(si, is_white)
    h_align_ok = paragraph_centered(si)
    v_align_ok = vertical_middle(si)

    passed = text_ok and centered_ok and font_ok and size_ok and bold_ok and color_ok and h_align_ok and v_align_ok
    detail = (
        f"文本={text_ok}；位于标题背景中央={centered_ok}（背景={'已找到' if bg else '未找到'}）；"
        f"宋体={font_ok}；12-14磅={size_ok}；加粗={bold_ok}；白色={color_ok}；"
        f"水平居中={h_align_ok}；垂直居中={v_align_ok}"
    )
    return CheckResult(5, "第1页顶部标题文本", passed, detail)


def body_frame_candidates(ctx: Context) -> list[ShapeInfo]:
    group = flow_chart_group_candidate(ctx)
    if not group:
        return []
    matched = []
    for si in ctx.shapes:
        if prst(si) != "roundRect":
            continue
        pos_ok = in_range(si.left - group.left, 0, 0.5, 0.0) and in_range(si.top - group.top, 0.6, 1.2, 0.0)
        size_ok = nearly(abs(si.width), 12.18, 0.01) and nearly(abs(si.height), 5.82, 0.01)
        fill_ok = is_white(fill_rgb(si))
        outline_ok = is_teal_or_green(line_rgb(si)) and is_solid_line(si)
        width_ok = line_width_ok(si, 0.75, 0.02)
        if pos_ok and size_ok and fill_ok and outline_ok and width_ok:
            matched.append(si)
    return matched


def check_body_frame(ctx: Context) -> CheckResult:
    group = flow_chart_group_candidate(ctx)
    if not group:
        return CheckResult(3, "第1页主体区域外框", False, "未找到完整流程图组合对象，无法定位主体区域外框")

    passed = bool(body_frame_candidates(ctx))
    detail = (
        "找到位于组合对象内距左0-0.5cm、距上0.6-1.2cm，12.18×5.82cm，白色圆角矩形，青绿色0.75磅实线轮廓的主体区域外框"
        if passed
        else "未找到同时满足位置、12.18×5.82cm、白色圆角矩形、青绿色实线轮廓、0.75磅线宽的主体区域外框"
    )
    return CheckResult(3, "第1页主体区域外框", passed, detail)


def shape_contains(outer: ShapeInfo, inner: ShapeInfo, tol: float = 0.01) -> bool:
    return (
        outer.left - tol <= inner.left
        and outer.top - tol <= inner.top
        and inner.right <= outer.right + tol
        and inner.bottom <= outer.bottom + tol
    )


def shape_center_inside(outer: ShapeInfo, inner: ShapeInfo, tol: float = 0.01) -> bool:
    cx, cy = shape_center(inner)
    return outer.left - tol <= cx <= outer.right + tol and outer.top - tol <= cy <= outer.bottom + tol


def check_background_bands(ctx: Context) -> CheckResult:
    body_frames = body_frame_candidates(ctx)
    if not body_frames:
        return CheckResult(5, "第1页主体区域背景分栏", False, "未找到主体外框，无法检查其内部背景分栏")
    body_frame = body_frames[0]

    bands = sorted(
        [
            si for si in ctx.shapes
            if prst(si) == "rect"
            and is_light_blue_gray(fill_rgb(si))
            and not norm_text(si.text)
            and shape_contains(body_frame, si, 0.01)
            and abs(si.height) > abs(si.width)
        ],
        key=lambda s: s.left,
    )

    count_ok = len(bands) == 10
    widths = [abs(b.width) for b in bands]
    tops = [b.top for b in bands]
    bottoms = [b.bottom for b in bands]
    equal_width = bool(widths) and max(widths) - min(widths) <= 0.01
    aligned = bool(tops) and max(tops) - min(tops) <= 0.01 and max(bottoms) - min(bottoms) <= 0.01

    boxes = ctx.step_boxes
    corresponding = False
    if len(boxes) == 10 and len(bands) == 10:
        corresponding = all(shape_center_inside(band, box, 0.01) for band, box in zip(bands, boxes))

    # 文字位置检查：文字可能写在步骤框自身，也可能放在独立文本框中——
    # 逐个步骤框调用 associated_text_shape 找到关联文本载体，再验证：
    #  - 若关联载体就是 box 自身（框内直接带字）→ 视为"文字位于框内"；
    #  - 若关联载体是独立文本框 → 要求其被步骤框合理包含
    #    （中心点落在 box 内 且 与 box 存在纵向交集，允许字体边距 0.1cm 容差）。
    def _text_carried_inside(box: ShapeInfo) -> bool:
        assoc = associated_text_shape(ctx, box)
        if assoc is None:
            return False
        if not norm_text(assoc.text):
            return False
        if assoc.index == box.index:
            return True
        # 独立文本框：中心落在 box 水平范围内 & 与 box 纵向有交集
        cx, cy = shape_center(assoc)
        if not (box.left - 0.10 <= cx <= box.right + 0.10):
            return False
        if not (box.top - 0.10 <= cy <= box.bottom + 0.10):
            return False
        # 纵向交集：assoc.top < box.bottom 且 assoc.bottom > box.top（容差 0.1cm）
        if assoc.bottom < box.top - 0.10 or assoc.top > box.bottom + 0.10:
            return False
        return True

    text_inside = len(boxes) == 10 and all(_text_carried_inside(box) for box in boxes)

    passed = count_ok and equal_width and aligned and corresponding and text_inside
    detail = (
        f"主体外框内浅蓝灰色竖向矩形色带数量={len(bands)}；数量为10={count_ok}；"
        f"宽度一致={equal_width}；顶部底部对齐={aligned}；"
        f"每个色带与一个处理步骤框对应={corresponding}；文字全部位于文本框内={text_inside}"
    )
    return CheckResult(5, "第1页主体区域背景分栏", passed, detail)


def check_step_boxes(ctx: Context) -> CheckResult:
    boxes = ctx.step_boxes

    # 细则：数量必须恰好为10
    count_ok = len(boxes) == 10
    if not count_ok:
        return CheckResult(5, "第1页第1-10个处理步骤框", False, f"仅识别到 {len(boxes)} 个处理步骤框，需恰好10个")

    # 细则：位于完整流程图组合对象内，距左0.40-12.30cm、距上1.5-4.5cm，宽0.80-0.90cm、高2-3cm
    # 十个步骤框分开独立识别，以组合对象为基准；无组合对象时直接检查各步骤框绝对尺寸
    group = flow_chart_group_candidate(ctx)
    origin_left = group.left if group else 0.0
    origin_top = group.top if group else 0.0
    pos_ok = all(
        in_range(b.left - origin_left, 0.40, 12.30, 0.0)
        and in_range(b.top - origin_top, 1.50, 4.50, 0.0)
        and in_range(abs(b.width), 0.80, 0.90, 0.0)
        and in_range(abs(b.height), 2.00, 3.00, 0.0)
        for b in boxes
    )

    # 细则：顶部编号分别为"1"-"10"（按从左往右顺序）
    # 数字文本可能位于 ellipse 内，也可能在紧邻的独立文本框内（例如两位数"10"）
    circles = ctx.number_circles
    nums_ok = False
    if len(circles) >= 10:
        def circle_number(c: ShapeInfo, expected: str) -> bool:
            if norm_text(c.text) == expected:
                return True
            # 兜底：在 ellipse 附近查找显示该数字的独立文本形状（水平/垂直交叠 ≤ 0.5cm）
            for s in ctx.shapes:
                if s.index == c.index:
                    continue
                if norm_text(s.text) != expected:
                    continue
                if s.right < c.left - 0.5 or s.left > c.right + 0.5:
                    continue
                if s.bottom < c.top - 0.5 or s.top > c.bottom + 0.5:
                    continue
                return True
            return False

        nums_ok = all(circle_number(circles[i], str(i + 1)) for i in range(10))

    # 细则：正文按从左往右顺序依次为10个指定步骤文本（文本可能在关联的独立文本框内）
    step_texts: list[str] = []
    for b in boxes:
        assoc = associated_text_shape(ctx, b)
        step_texts.append(assoc.text if assoc is not None else b.text)
    texts_ok = all(expected_step_text_ok(i, t) for i, t in enumerate(step_texts))

    # 细则：每两个处理步骤框的间距为0.3-0.5cm
    spacing = [boxes[i + 1].left - boxes[i].right for i in range(9)]
    spacing_ok = all(in_range(s, 0.30, 0.50, 0.0) for s in spacing)

    passed = pos_ok and nums_ok and texts_ok and spacing_ok
    detail = (
        f"数量=10={count_ok}；位置尺寸={pos_ok}（基准={'组合对象' if group else '幻灯片原点'}）；"
        f"顶部编号1-10={nums_ok}；正文顺序={texts_ok}；"
        f"间距0.3-0.5cm={spacing_ok}（各间距={[round(s, 2) for s in spacing]}）"
    )
    return CheckResult(5, "第1页第1-10个处理步骤框", passed, detail)


def check_step_box_style(ctx: Context) -> CheckResult:
    boxes = ctx.step_boxes
    if len(boxes) != 10:
        return CheckResult(5, "第1页处理步骤框样式", False, f"仅识别到 {len(boxes)} 个步骤框，需恰好10个")

    # 细则：白色圆角矩形
    shape_ok = all(prst(b) == "roundRect" and is_white_or_no_fill(b) for b in boxes)
    # 细则：轮廓为青色或绿色实线
    outline_ok = all(is_teal_or_green(line_rgb(b)) and is_solid_line(b) for b in boxes)
    # 细则：线宽0.75磅
    linewidth_ok = all(line_width_ok(b, 0.75, 0.02) for b in boxes)

    # 细则：顶部和底部对齐
    tops = [b.top for b in boxes]
    bottoms = [b.bottom for b in boxes]
    align_ok = max(tops) - min(tops) <= 0.01 and max(bottoms) - min(bottoms) <= 0.01

    # 细则：宽度保持一致 & 相邻框间距保持一致
    # 容差 0.05cm ≈ 半个字号点距，可容纳浮点/取整误差，但显著不一致会被判失败。
    widths = [abs(b.width) for b in boxes]
    width_uniform = (max(widths) - min(widths)) <= 0.05
    boxes_by_left = sorted(boxes, key=lambda s: s.left)
    spacing = [boxes_by_left[i + 1].left - boxes_by_left[i].right for i in range(9)]
    spacing_uniform = bool(spacing) and (max(spacing) - min(spacing)) <= 0.05

    passed = shape_ok and outline_ok and linewidth_ok and align_ok and width_uniform and spacing_uniform
    detail = (
        f"白色圆角矩形={shape_ok}；青绿实线轮廓={outline_ok}；0.75磅线宽={linewidth_ok}；"
        f"顶部底部对齐={align_ok}；"
        f"宽度一致={width_uniform}(极差={round(max(widths) - min(widths), 3)}cm)；"
        f"间距一致={spacing_uniform}(极差={round(max(spacing) - min(spacing), 3) if spacing else 0}cm)"
    )
    return CheckResult(5, "第1页处理步骤框样式", passed, detail)


def step_text_is_vertical_layout(b: ShapeInfo) -> bool:
    """逐字换行：文本框内每行不超过1个汉字（即文本中有换行），或文本框宽度<=0.90cm（窄文本框）。"""
    text = b.text
    has_newline = "\n" in text
    is_narrow = abs(b.width) <= 0.90
    return has_newline or is_narrow


def check_step_text_style(ctx: Context) -> CheckResult:
    boxes = ctx.step_boxes
    if len(boxes) != 10:
        return CheckResult(5, "第1页处理步骤正文", False, f"仅识别到 {len(boxes)} 个步骤框")

    # 找到每个步骤框对应的正文形状（可能是关联的独立文本框，也可能是框自身）
    text_shapes: list[ShapeInfo] = []
    for b in boxes:
        assoc = associated_text_shape(ctx, b)
        text_shapes.append(assoc if assoc is not None else b)

    # 细则：宋体
    font_ok = all(has_song_font(t) for t in text_shapes)
    # 细则：5-8磅
    size_ok = all(sizes_in(t, 5, 8) for t in text_shapes)
    # 细则：加粗
    bold_ok = all(all_bold(t) for t in text_shapes)
    # 细则：黑色
    color_ok = all(text_colors_match(t, is_black) for t in text_shapes)
    # 细则：逐字换行或窄文本框排列
    layout_ok = all(step_text_is_vertical_layout(t) for t in text_shapes)
    # 细则：水平居中
    h_align_ok = all(paragraph_centered(t) for t in text_shapes)
    # 细则：垂直居中
    v_align_ok = all(vertical_middle(t) for t in text_shapes)

    passed = font_ok and size_ok and bold_ok and color_ok and layout_ok and h_align_ok and v_align_ok
    detail = (
        f"宋体={font_ok}；5-8磅={size_ok}；加粗={bold_ok}；黑色={color_ok}；"
        f"逐字换行或窄框={layout_ok}；水平居中={h_align_ok}；垂直居中={v_align_ok}"
    )
    return CheckResult(5, "第1页处理步骤正文", passed, detail)


def number_circle_at_box_top(circle: ShapeInfo, box: ShapeInfo) -> bool:
    return (
        box.left <= circle.left
        and circle.right <= box.right
        and box.top <= circle.top
        and circle.bottom <= box.bottom
        and circle.top <= box.top + 0.50
    )


def check_number_circles(ctx: Context) -> CheckResult:
    circles = ctx.number_circles
    if len(circles) != 10:
        return CheckResult(5, "第1页处理步骤编号圆形", False, f"识别到 {len(circles)} 个编号圆形，需恰好10个")

    # 圆内没有文字时（例如两位数"10"被放到圆外的独立文本框），到旁边找一个显示纯数字的文本形状，
    # 用它承担字体/字号/加粗/颜色的判定
    def number_text_shape(c: ShapeInfo) -> ShapeInfo:
        if norm_text(c.text):
            return c
        for s in ctx.shapes:
            if s.index == c.index:
                continue
            t = norm_text(s.text)
            if not (t and t.isdigit()):
                continue
            if s.right < c.left - 0.5 or s.left > c.right + 0.5:
                continue
            if s.bottom < c.top - 0.5 or s.top > c.bottom + 0.5:
                continue
            return s
        return c

    text_shapes = [number_text_shape(c) for c in circles]

    # 细则：10个编号圆形直径均为0.25-0.30厘米
    diameter_ok = all(
        prst(c) == "ellipse"
        and in_range(abs(c.width), 0.25, 0.30, 0.0)
        and in_range(abs(c.height), 0.25, 0.30, 0.0)
        and abs(abs(c.width) - abs(c.height)) <= 0.01
        for c in circles
    )
    # 细则：填充为青色或绿色
    fill_ok = all(is_teal_or_green(fill_rgb(c)) for c in circles)
    # 细则：轮廓为青绿色
    outline_ok = all(is_teal_or_green(line_rgb(c)) for c in circles)
    # 细则：轮廓线宽0.75磅
    circle_linewidth_ok = all(line_width_ok(c, 0.75, 0.02) for c in circles)
    # 细则：数字为Times New Roman9磅、加粗、白色
    font_ok = all(has_times_font(t) for t in text_shapes)
    size_ok = all(sizes_in(t, 9, 9, 0.1) for t in text_shapes)
    bold_ok = all(all_bold(t) for t in text_shapes)
    color_ok = all(text_colors_match(t, is_white) for t in text_shapes)
    # 细则：放置在处理步骤框内顶部
    top_inside_ok = len(ctx.step_boxes) == 10 and all(number_circle_at_box_top(c, b) for c, b in zip(circles, ctx.step_boxes))

    passed = (
        diameter_ok and fill_ok and outline_ok and circle_linewidth_ok
        and font_ok and size_ok and bold_ok and color_ok and top_inside_ok
    )
    detail = (
        f"数量=10；直径0.25-0.30cm={diameter_ok}；填充青/绿={fill_ok}；"
        f"轮廓青绿色={outline_ok}；轮廓0.75磅={circle_linewidth_ok}；"
        f"Times New Roman={font_ok}；"
        f"9磅={size_ok}；加粗={bold_ok}；白色={color_ok}；位于步骤框内顶部={top_inside_ok}"
    )
    return CheckResult(5, "第1页处理步骤编号圆形", passed, detail)


def arrow_between_step_boxes(arrow: ShapeInfo, prev_box: ShapeInfo, next_box: ShapeInfo) -> bool:
    (x1, y1), (x2, y2) = line_points(arrow)
    # 步骤间连接可能是普通直线 (line) 或直线箭头连接符 (straightConnector1)
    if not is_line_like(arrow):
        return False
    # 颜色：优先按明确指定的 RGB 判断；若使用主题色（line_rgb 返回 None），只要走 lnRef 主题线样式即认可
    color_rgb = line_rgb(arrow)
    color_ok = is_dark_blue(color_rgb) or (color_rgb is None and _uses_line_style_ref(arrow))
    if not color_ok:
        return False
    return (
        is_horizontal(arrow, 0.01)
        and has_any_arrow(arrow)
        and points_left_to_right(arrow)
        and line_width_ok(arrow, 0.75, 0.02)
        and in_range(abs(arrow.width), 0.25, 0.50, 0.0)
        and abs(x1 - prev_box.right) <= 0.03
        and abs(x2 - next_box.left) <= 0.03
        and prev_box.top <= y1 <= prev_box.bottom
        and next_box.top <= y2 <= next_box.bottom
    )


def check_horizontal_arrows(ctx: Context) -> CheckResult:
    if len(ctx.step_boxes) != 10:
        return CheckResult(5, "第1页步骤间水平箭头", False, "未能识别完整10个步骤框，无法验证箭头")

    line_candidates = [s for s in ctx.shapes if is_line_like(s)]
    matched = []
    for prev_box, next_box in zip(ctx.step_boxes, ctx.step_boxes[1:]):
        pair_matches = [s for s in line_candidates if arrow_between_step_boxes(s, prev_box, next_box)]
        matched.append(len(pair_matches) == 1)

    passed = all(matched)
    detail = (
        f"9处步骤框间均有且仅有1条深蓝色水平单箭头={passed}；"
        f"逐项={matched}；要求：从左到右、前框右边缘到后框左边缘、0.75磅、长度0.25-0.5cm"
    )
    return CheckResult(5, "第1页步骤间水平箭头", passed, detail)



def icon_near_left(ctx: Context, box: ShapeInfo) -> bool:
    # 左侧图标可能在文本框内部拆分为独立可编辑矢量形状，也可能贴在左边外侧。
    candidates = []
    for si in ctx.shapes:
        if si.index == box.index:
            continue
        if not is_vector_icon_candidate(si):
            continue
        center_x = si.left + si.width / 2
        center_y = si.top + si.height / 2
        if box.left - 0.25 <= center_x <= box.left + 0.35 and box.top - 0.15 <= center_y <= box.bottom + 0.15:
            candidates.append(si)
    return bool(candidates)


def _aux_box_check(ctx: Context, si: ShapeInfo, points: int, name: str,
                   left_rng: tuple, top_rng: tuple, w_rng: tuple, h_rng: tuple,
                   icon_name: str) -> CheckResult:
    group = flow_chart_group_candidate(ctx)
    if not group:
        return CheckResult(points, f"第1页{name}说明框", False, "未找到完整流程图组合对象，无法验证位置")
    text_ok = norm_text(si.text) == name
    left_ok = in_range(si.left - group.left, left_rng[0], left_rng[1], 0.0)
    top_ok = in_range(si.top - group.top, top_rng[0], top_rng[1], 0.0)
    width_ok = in_range(abs(si.width), w_rng[0], w_rng[1], 0.0)
    height_ok = in_range(abs(si.height), h_rng[0], h_rng[1], 0.0)
    icon_ok = icon_near_left(ctx, si)
    passed = text_ok and left_ok and top_ok and width_ok and height_ok and icon_ok
    detail = (
        f"文本为{name}={text_ok}；"
        f"距组合对象左{left_rng[0]}-{left_rng[1]}cm={left_ok}（实际{si.left - group.left:.2f}cm）；"
        f"距组合对象上{top_rng[0]}-{top_rng[1]}cm={top_ok}（实际{si.top - group.top:.2f}cm）；"
        f"宽{w_rng[0]}-{w_rng[1]}cm={width_ok}（实际{abs(si.width):.2f}cm）；"
        f"高{h_rng[0]}-{h_rng[1]}cm={height_ok}（实际{abs(si.height):.2f}cm）；"
        f"左侧可编辑{icon_name}图标={icon_ok}"
    )
    return CheckResult(points, f"第1页{name}说明框", passed, detail)


def check_aux_box(spec) -> Callable[[Context], CheckResult]:
    text, points, x_rng, y_rng, w_rng, h_rng, icon_name = spec

    def _check(ctx: Context) -> CheckResult:
        si = find_text_shape(ctx, text)
        if not si:
            return CheckResult(points, f"第1页{text}说明框", False, f"未找到文本【{text}】")

        if text == "酸碱调节":
            return _aux_box_check(ctx, si, points, "酸碱调节",
                                  (2.5, 4.0), (4.3, 5.0), (1.0, 1.3), (0.3, 0.5), "烧瓶")
        if text == "PAC/PAM投加":
            return _aux_box_check(ctx, si, points, "PAC/PAM投加",
                                  (3.5, 5.0), (4.3, 5.0), (1.3, 1.4), (0.3, 0.4), "药剂")
        if text == "核心净化单元":
            return _aux_box_check(ctx, si, points, "核心净化单元",
                                  (6.4, 8.4), (4.3, 5.0), (1.4, 1.6), (0.3, 0.4), "水滴")
        if text == "次氯酸钠投加":
            return _aux_box_check(ctx, si, points, "次氯酸钠投加",
                                  (8.3, 10.1), (4.3, 5.0), (1.4, 1.6), (0.3, 0.4), "盾牌")

        # 其余说明框沿用旧逻辑
        cl, ct, _, _ = ctx.chart
        pos_ok = rect_close(si, x_rng, y_rng, w_rng, h_rng, (cl, ct), 0.08)
        icon_ok = icon_near_left(ctx, si)
        passed = pos_ok and icon_ok
        detail = f"位置尺寸={pos_ok}；左侧可编辑{icon_name}图标={icon_ok}；实际 local=({si.left-cl:.2f},{si.top-ct:.2f},{abs(si.width):.2f},{abs(si.height):.2f})"
        return CheckResult(points, f"第1页{text}说明框", passed, detail)

    return _check


def is_deep_green_text(c: Optional[str]) -> bool:
    return is_dark_teal(c) or color_near(c, ["006B60", "00796B", "008060", "006400", "008000"], 55)


def check_aux_style(ctx: Context) -> CheckResult:
    boxes = [find_text_shape(ctx, spec[0]) for spec in AUX_SPECS]
    if not all(boxes):
        return CheckResult(3, "第1页辅助说明框样式", False, f"4个说明框未全部找到：{[bool(b) for b in boxes]}")
    assert all(boxes)

    # 细则：4个说明框均为极浅绿色圆角矩形
    shape_fill_ok = all(prst(b) == "roundRect" and is_light_green(fill_rgb(b)) for b in boxes if b)
    # 细则：轮廓为青绿色0.75磅
    outline_ok = all(is_teal_or_green(line_rgb(b)) and line_width_ok(b, 0.75, 0.02) for b in boxes if b)
    # 细则：文字为宋体，大小5-7磅、深绿色
    font_ok = all(has_song_font(b) for b in boxes if b)
    size_ok = all(sizes_in(b, 5, 7, 0.1) for b in boxes if b)
    color_ok = all(text_colors_match(b, is_deep_green_text) for b in boxes if b)

    passed = shape_fill_ok and outline_ok and font_ok and size_ok and color_ok
    detail = (
        f"4个说明框均为极浅绿色圆角矩形={shape_fill_ok}；"
        f"青绿色0.75磅轮廓={outline_ok}；宋体={font_ok}；"
        f"5-7磅={size_ok}；深绿色文字={color_ok}"
    )
    return CheckResult(3, "第1页辅助说明框样式", passed, detail)


def check_aux_dashed_arrows(ctx: Context) -> CheckResult:
    if len(ctx.step_boxes) < 8:
        return CheckResult(3, "第1页辅助说明虚线箭头", False, "未识别到足够步骤框")
    # 细则：连接第3、第4、第7、第8个步骤框（索引2,3,6,7）
    target_indices = [2, 3, 6, 7]
    aux_boxes = [find_text_shape(ctx, spec[0]) for spec in AUX_SPECS]

    # 候选线段：青绿色、竖向、虚线（同时接受 line 与 straightConnector1）
    segs = [
        s for s in ctx.shapes
        if is_line_like(s)
        and is_teal_or_green(line_rgb(s))
        and is_vertical(s, 0.01)
        and is_dashed(s)
    ]

    item_ok = []
    for step_idx, aux in zip(target_indices, aux_boxes):
        if not aux:
            item_ok.append(False)
            continue
        step = ctx.step_boxes[step_idx]
        found = False
        for s in segs:
            # 细则：线宽0.75磅
            if not line_width_ok(s, 0.75, 0.02):
                continue
            # 细则：长度0.3-0.5cm
            length = abs(s.height)
            if not in_range(length, 0.30, 0.50, 0.0):
                continue
            # 端点：aux 框位于 step 下方，故线段 y_low≈step.bottom，y_high≈aux.top
            y_low, y_high = sorted([s.top, s.top + s.height])
            touches_step_bottom = abs(y_low - step.bottom) <= 0.05
            touches_aux_top = abs(y_high - aux.top) <= 0.05
            if not (touches_step_bottom and touches_aux_top):
                continue
            # x 对齐：线段 x 落在 step 或 aux 的中心附近（容差 0.20cm）
            sx = s.left
            step_cx = step.left + abs(step.width) / 2
            aux_cx = aux.left + abs(aux.width) / 2
            x_aligned = abs(sx - step_cx) <= 0.20 or abs(sx - aux_cx) <= 0.20
            if not x_aligned:
                continue
            # 细则：箭头方向向上——准确解析 headEnd/tailEnd 与端点方向：
            #   line_points 已尊重 flipV，(x1,y1) 为 tail、(x2,y2) 为 head；
            #   若存在 headEnd，箭头位于 head 端，需 y2 < y1（箭头端更靠上）；
            #   若仅存在 tailEnd，箭头位于 tail 端，需 y1 < y2；
            #   两端都无箭头 → 不合格。
            (_x1, _y1), (_x2, _y2) = line_points(s)
            if has_head_arrow(s):
                arrow_up = _y2 < _y1
            elif has_tail_arrow(s):
                arrow_up = _y1 < _y2
            else:
                arrow_up = False
            if not arrow_up:
                continue
            found = True
            break
        item_ok.append(found)

    passed = all(item_ok)
    detail = (
        f"4条辅助说明虚线箭头全部通过={passed}；逐项={item_ok}；"
        f"要求：青绿色竖向虚线、0.75磅、长0.3-0.5cm、向上箭头、连接说明框与对应步骤框底部"
    )
    return CheckResult(3, "第1页辅助说明虚线箭头", passed, detail)


def check_return_pipeline(ctx: Context) -> CheckResult:
    if len(ctx.step_boxes) < 9:
        return CheckResult(5, "第1页异常回流管线", False, "未识别到第2个或第9个步骤框")

    step2 = ctx.step_boxes[1]
    step9 = ctx.step_boxes[8]
    body_frames = body_frame_candidates(ctx)
    body_frame = body_frames[0] if body_frames else None

    horizontals = [
        s for s in ctx.shapes
        if is_line_like(s)
        and (is_dark_blue(line_rgb(s)) or (line_rgb(s) is None and _uses_line_style_ref(s)))
        and is_horizontal(s, 0.05)
        and line_width_ok(s, 0.75, 0.05)
    ]
    verticals = [
        s for s in ctx.shapes
        if is_line_like(s)
        and (is_dark_blue(line_rgb(s)) or (line_rgb(s) is None and _uses_line_style_ref(s)))
        and is_vertical(s, 0.05)
        and line_width_ok(s, 0.75, 0.05)
        and in_range(abs(s.height), 1.0, 1.5, 0.05)
        and has_any_arrow(s)
        and points_up(s)
    ]

    # 细则要求水平段长约 7.7-7.9cm；实际 PPT 可能被“异常回流”标签分隔为多段，
    # 因此以“共同高度带内所有段的合并左右跨度”作为水平总长度评估。
    matched = False
    matched_h_group_len = 0.0
    for v in verticals:
        vx = v.left
        vy1, vy2 = sorted([v.top, v.top + v.height])
        # 竖直向上箭头应位于第2个步骤框下方并连接其底部
        v_under_step2 = step2.left - 0.10 <= vx <= step2.right + 0.10
        v_connects_step2 = abs(vy1 - step2.bottom) <= 0.10 or (step2.top <= vy1 <= step2.bottom + 0.10)
        if not (v_under_step2 and v_connects_step2):
            continue
        # 从竖向顶点(vx, vy2)出发的水平链：找与之处于相似 Y 高度的所有水平段
        chain = [h for h in horizontals if abs(h.top - vy2) <= 0.30]
        if not chain:
            continue
        left_min = min(h.left for h in chain)
        right_max = max(h.left + h.width for h in chain)
        # 合并跨度：从第2步下方竖向 x 到第9步下方
        total_span = right_max - left_min
        step9_reach = right_max >= step9.left - 0.10
        vx_covered = abs(left_min - vx) <= 0.30 or left_min <= vx + 0.30
        # body 区域下部
        h_in_body_lower = True if not body_frame else all(
            body_frame.top + abs(body_frame.height) * 0.55 <= h.top <= body_frame.bottom + 0.20
            for h in chain
        )
        if step9_reach and vx_covered and h_in_body_lower and in_range(total_span, 7.5, 8.0, 0.10):
            matched = True
            matched_h_group_len = total_span
            break

    h_ok = bool(horizontals)
    v_ok = bool(verticals)
    passed = h_ok and v_ok and matched
    detail = (
        f"深蓝色0.75磅水平段(可多段){h_ok}；"
        f"深蓝色0.75磅竖直向上箭头且长1.0-1.5cm={v_ok}；"
        f"合并水平跨度覆盖第9→第2步={matched}（合并跨度{matched_h_group_len:.2f}cm）"
    )
    return CheckResult(5, "第1页异常回流管线", passed, detail)


def check_return_label(ctx: Context) -> CheckResult:
    si = find_text_shape(ctx, "异常回流")
    if not si:
        return CheckResult(3, "第1页异常回流标签", False, "未找到文本【异常回流】")
    group = flow_chart_group_candidate(ctx)
    if not group:
        return CheckResult(3, "第1页异常回流标签", False, "未找到完整流程图组合对象，无法验证位置")

    # 细则：位于完整流程图组合对象内距左4.8-6cm、距上5.3-5.8cm，宽0.9-1.1cm、高0.3-0.4cm
    left_ok = in_range(si.left - group.left, 4.8, 6.0, 0.0)
    top_ok = in_range(si.top - group.top, 5.3, 5.8, 0.0)
    width_ok = in_range(abs(si.width), 0.9, 1.1, 0.0)
    height_ok = in_range(abs(si.height), 0.3, 0.4, 0.0)
    # 细则：填充为深蓝色（PPT 常见做法：不显式 solidFill，通过 p:style/fillRef 指向主题 accent1 → 4874CB）
    fill_ok = is_module_main_blue(fill_rgb_with_style(si))
    # 细则：文本为“异常回流”
    text_ok = norm_text(si.text) == "异常回流"
    # 细则：文字为宋体，大小5-7磅、加粗、白色
    font_ok = has_song_font(si)
    size_ok = sizes_in(si, 5, 7, 0.1)
    bold_ok = all_bold(si)
    color_ok = text_colors_match(si, is_white)

    passed = left_ok and top_ok and width_ok and height_ok and fill_ok and text_ok and font_ok and size_ok and bold_ok and color_ok
    detail = (
        f"距组合对象左4.8-6cm={left_ok}（实际{si.left - group.left:.2f}cm）；"
        f"距组合对象上5.3-5.8cm={top_ok}（实际{si.top - group.top:.2f}cm）；"
        f"宽0.9-1.1cm={width_ok}（实际{abs(si.width):.2f}cm）；"
        f"高0.3-0.4cm={height_ok}（实际{abs(si.height):.2f}cm）；"
        f"深蓝填充={fill_ok}；文本异常回流={text_ok}；宋体={font_ok}；"
        f"5-7磅={size_ok}；加粗={bold_ok}；白色={color_ok}"
    )
    return CheckResult(3, "第1页异常回流标签", passed, detail)


def check_reuse_box(ctx: Context) -> CheckResult:
    si = find_text_shape(ctx, "绿化与冲洗回用")
    if not si:
        return CheckResult(5, "第1页绿化与冲洗回用框", False, "未找到文本【绿化与冲洗回用】")
    group = flow_chart_group_candidate(ctx)
    if not group:
        return CheckResult(5, "第1页绿化与冲洗回用框", False, "未找到完整流程图组合对象，无法验证位置")

    # 细则：位于完整流程图组合对象内距左9.6-12.3cm、距上5.10-5.8cm，宽2.1-2.3cm、高0.5-0.6cm
    left_ok = in_range(si.left - group.left, 9.6, 12.3, 0.0)
    top_ok = in_range(si.top - group.top, 5.10, 5.8, 0.0)
    width_ok = in_range(abs(si.width), 2.1, 2.3, 0.0)
    height_ok = in_range(abs(si.height), 0.5, 0.6, 0.0)
    # 细则：形状为白色圆角矩形
    shape_ok = prst(si) == "roundRect" and is_white(fill_rgb(si))
    # 细则：轮廓为青绿色0.75磅
    outline_ok = is_teal_or_green(line_rgb(si)) and line_width_ok(si, 0.75, 0.02)
    # 细则：左侧为绿色叶片图标（可编辑矢量图标）
    icon_ok = icon_near_left(ctx, si)
    # 细则：文本为“绿化与冲洗回用”
    text_ok = norm_text(si.text) == "绿化与冲洗回用"

    passed = left_ok and top_ok and width_ok and height_ok and shape_ok and outline_ok and icon_ok and text_ok
    detail = (
        f"距组合对象左9.6-12.3cm={left_ok}（实际{si.left - group.left:.2f}cm）；"
        f"距组合对象上5.10-5.8cm={top_ok}（实际{si.top - group.top:.2f}cm）；"
        f"宽2.1-2.3cm={width_ok}（实际{abs(si.width):.2f}cm）；"        f"高0.5-0.6cm={height_ok}（实际{abs(si.height):.2f}cm）；"
        f"白色圆角矩形={shape_ok}；青绿色0.75磅轮廓={outline_ok}；"
        f"左侧绿色叶片图标={icon_ok}；文本绿化与冲洗回用={text_ok}"
    )
    return CheckResult(5, "第1页绿化与冲洗回用框", passed, detail)


def check_down_arrow(ctx: Context) -> CheckResult:
    target = find_text_shape(ctx, "绿化与冲洗回用")
    if not target:
        return CheckResult(3, "第1页回用清水池下行箭头", False, "未找到【绿化与冲洗回用】框")
    if len(ctx.step_boxes) < 9:
        return CheckResult(3, "第1页回用清水池下行箭头", False, "未识别到第9个步骤框")
    step9 = ctx.step_boxes[8]

    candidates = [
        s for s in ctx.shapes
        if is_line_like(s)
        and (is_dark_blue(line_rgb(s)) or (line_rgb(s) is None and _uses_line_style_ref(s)))
        and is_vertical(s, 0.05)
        and line_width_ok(s, 0.75, 0.05)
        and in_range(abs(s.height), 1.0, 1.2, 0.05)
        and has_any_arrow(s)
        and points_down(s)
    ]

    found = False
    for s in candidates:
        x = s.left
        y1, y2 = sorted([s.top, s.top + s.height])
        # 细则：从第9个步骤框底部出发，向下连接绿化与冲洗回用框顶部
        from_step9_bottom = abs(y1 - step9.bottom) <= 0.10
        to_target_top = abs(y2 - target.top) <= 0.10
        x_aligned = (step9.left - 0.10 <= x <= step9.right + 0.10) and (target.left - 0.10 <= x <= target.right + 0.10)
        if from_step9_bottom and to_target_top and x_aligned:
            found = True
            break

    detail = (
        "找到从第9步框底部垂直向下、连接绿化与冲洗回用框顶部的深蓝色0.75磅单向下行箭头（长1-1.2cm）"
        if found
        else (
            f"未找到满足全部条件的下行箭头；候选深蓝色竖向箭头数={len(candidates)}；"
            f"第9步底部y={step9.bottom:.2f}cm；绿化框顶部y={target.top:.2f}cm"
        )
    )
    return CheckResult(3, "第1页回用清水池下行箭头", found, detail)



def check_bottom_bar(ctx: Context) -> CheckResult:
    group = flow_chart_group_candidate(ctx)
    if not group:
        return CheckResult(5, "第1页底部保障信息栏", False, "未找到完整流程图组合对象，无法验证底部保障信息栏位置")

    matched = [
        si for si in ctx.shapes
        if prst(si) == "roundRect"
        and in_range(si.left - group.left, 0.3, 12.2, 0.05)
        and in_range(si.top - group.top, 5.8, 6.8, 0.05)
        and in_range(abs(si.width), 11.0, 12.0, 0.10)
        and in_range(abs(si.height), 0.6, 0.8, 0.10)
        and is_light_blue(fill_rgb(si))
        and is_teal_or_green(line_rgb(si))
        and line_width_ok(si, 0.75, 0.05)
    ]
    if matched:
        si = matched[0]
        detail = (
            f"找到极浅蓝色圆角矩形底部信息栏；"
            f"距组合对象左{si.left - group.left:.2f}cm；距组合对象上{si.top - group.top:.2f}cm；"
            f"宽{abs(si.width):.2f}cm；高{abs(si.height):.2f}cm；蓝绿色0.75磅轮廓"
        )
        return CheckResult(5, "第1页底部保障信息栏", True, detail)
    return CheckResult(5, "第1页底部保障信息栏", False,
                       "未找到同时满足位置（距左0.3-12.2cm、距上5.8-6.8cm）、尺寸（11-12×0.6-0.8cm）、圆角矩形、极浅蓝色填充、蓝绿色0.75磅轮廓的底部保障信息栏")


def check_bottom_dividers(ctx: Context) -> CheckResult:
    group = flow_chart_group_candidate(ctx)
    if not group:
        return CheckResult(3, "第1页底部保障信息栏分隔线", False, "未找到完整流程图组合对象，无法定位底部信息栏")

    # 找底部信息栏（复用 check_bottom_bar 的定位逻辑）
    bar_candidates = [
        si for si in ctx.shapes
        if prst(si) == "roundRect"
        and in_range(si.left - group.left, 0.3, 12.2, 0.05)
        and in_range(si.top - group.top, 5.8, 6.8, 0.05)
        and in_range(abs(si.width), 11.0, 12.0, 0.10)
        and in_range(abs(si.height), 0.6, 0.8, 0.10)
        and is_teal_or_green(line_rgb(si))
        and line_width_ok(si, 0.75, 0.05)
    ]
    bar = bar_candidates[0] if bar_candidates else None

    # 候选线段：深蓝或蓝绿色、竖向、虚线、位于信息栏内部（线宽放宽以兼容 1pt 实际值）
    segs = [
        s for s in ctx.shapes
        if is_line_like(s)
        and (is_teal_or_green(line_rgb(s)) or is_dark_blue(line_rgb(s))
             or (line_rgb(s) is None and _uses_line_style_ref(s)))
        and is_vertical(s, 0.05)
        and is_dashed(s)
        and line_width_ok(s, 0.75, 0.30)
        and in_range(abs(s.height), 0.4, 0.6, 0.05)
        and (bar is None or (bar.left - 0.05 <= s.left <= bar.right + 0.05
                              and bar.top - 0.05 <= s.top
                              and s.top + abs(s.height) <= bar.bottom + 0.05))
    ]

    # 细则：3条虚线将信息栏划分为4个等宽区域（考虑圆角矩形两端受曲线影响，容差放宽到 0.5cm）
    count_ok = len(segs) == 3
    equal_spacing = False
    spacings = []
    if count_ok and bar is not None:
        xs = sorted(s.left for s in segs)
        spacings = [xs[0] - bar.left, xs[1] - xs[0], xs[2] - xs[1], bar.right - xs[2]]
        equal_spacing = max(spacings) - min(spacings) <= 0.50
    elif count_ok:
        equal_spacing = True

    passed = count_ok and equal_spacing
    detail = (
        f"蓝绿色竖向虚线0.75磅且长0.4-0.6cm数量={len(segs)}（需3条）；"
        f"将信息栏划分为4个等宽区域={equal_spacing}"
        + (f"；各区域宽={[round(s, 2) for s in spacings]}" if count_ok and bar else "")
    )
    return CheckResult(3, "第1页底部保障信息栏分隔线", passed, detail)


def check_bottom_module(index: int, main: str, sub: str, icon_name: str) -> Callable[[Context], CheckResult]:
    def _check(ctx: Context) -> CheckResult:
        main_shape, sub_shape = find_module_pair(ctx, main, sub)
        if not main_shape and not sub_shape:
            return CheckResult(1, f"第1页底部第{index}个保障模块", False, f"未找到主文本【{main}】和副文本【{sub}】")

        main_ok = main_shape is not None
        sub_ok = sub_shape is not None
        # 图标位于主文本框左侧（若拆分为两个 rect，则以 main 为基准；否则以合并框为基准）
        anchor = main_shape or sub_shape
        assert anchor is not None
        icon_ok = icon_near_left(ctx, anchor)
        passed = icon_ok and main_ok and sub_ok
        detail = (
            f"可编辑{icon_name}图标={icon_ok}；"
            f"主文本{main}={main_ok}；"
            f"副文本{sub}={sub_ok}"
        )
        return CheckResult(1, f"第1页底部第{index}个保障模块", passed, detail)
    return _check


def runs_match_style(runs, *, size_target: float, bold: Optional[bool], color_pred: Callable[[Optional[str]], bool]) -> bool:
    if not runs:
        return False
    names = [run.font.name or "" for run in runs]
    song_ok = (not names) or any("宋" in n or "SimSun" in n or "Song" in n for n in names)
    sizes = [run.font.size.pt for run in runs if run.font.size]
    size_ok = bool(sizes) and all(abs(s - size_target) <= 0.1 for s in sizes)
    colors = [safe_rgb(run.font.color) for run in runs]
    color_ok = all(color_pred(c) for c in colors if c is not None)
    bold_ok = True if bold is None else all(run.font.bold is bold for run in runs)
    return song_ok and size_ok and color_ok and bold_ok


def check_bottom_text_style(ctx: Context) -> CheckResult:
    pairs = [find_module_pair(ctx, main, sub) for main, sub, _ in BOTTOM_MODULES]
    if not all(m or s for m, s in pairs):
        return CheckResult(1, "第1页底部保障模块文字", False,
                           f"4个模块文本未全部找到：{[(bool(m), bool(s)) for m, s in pairs]}")

    ok = True
    detail_parts = []

    def _runs_of(shape: ShapeInfo, needle: str) -> list:
        try:
            for p in shape.shape.text_frame.paragraphs:
                ptxt = norm_text("".join(run.text for run in p.runs)).lower()
                if needle.lower() in ptxt:
                    return [r for r in p.runs if r.text]
        except Exception:
            return []
        return []

    for (main_shape, sub_shape), (main, sub, _) in zip(pairs, BOTTOM_MODULES):
        main_ok = False
        sub_ok = False
        if main_shape is not None:
            runs = _runs_of(main_shape, main)
            # 细则：主文本使用宋体6磅、加粗、深蓝色
            if runs:
                main_ok = runs_match_style(runs, size_target=6, bold=True, color_pred=is_module_main_blue)
        if sub_shape is not None:
            runs = _runs_of(sub_shape, sub)
            # 细则：副文本为灰色，使用宋体五磅
            if runs:
                sub_ok = runs_match_style(runs, size_target=5, bold=None, color_pred=is_gray)
        detail_parts.append(f"{main}:主文本宋体6磅加粗深蓝={main_ok},副文本宋体5磅灰色={sub_ok}")
        ok = ok and main_ok and sub_ok
    return CheckResult(1, "第1页底部保障模块文字", ok, "；".join(detail_parts))


CHECKS: list[Callable[[Context], CheckResult]] = [
    check_complete_group,
    check_title_bg,
    check_title_text,
    check_body_frame,
    check_background_bands,
    check_step_boxes,
    check_step_box_style,
    check_step_text_style,
    check_number_circles,
    check_horizontal_arrows,
    *[check_aux_box(spec) for spec in AUX_SPECS],
    check_aux_style,
    check_aux_dashed_arrows,
    check_return_pipeline,
    check_return_label,
    check_reuse_box,
    check_down_arrow,
    check_bottom_bar,
    check_bottom_dividers,
    *[check_bottom_module(i + 1, main, sub, icon) for i, (main, sub, icon) in enumerate(BOTTOM_MODULES)],
    check_bottom_text_style,
]


SCRIPT_ID = "054"
MAX_SCORE = 89  # 维度二所有评分项 max_delta 之和


def _locate_target(dir_path: Path) -> Optional[Path]:
    """在 dir_path 目录里定位被评估文档：优先精确匹配 TARGET_FILE，否则取第一个 .pptx。"""
    target = dir_path / TARGET_FILE
    if target.exists():
        return target
    try:
        candidates = sorted(
            p for p in dir_path.iterdir()
            if p.is_file() and p.suffix.lower() == ".pptx"
        )
    except Exception:
        return None
    return candidates[0] if candidates else None


def evaluate(dir_path: str) -> dict:
    """在 dir_path 目录里定位并评估 PPT 文档，返回结构化评估结果。"""
    result: dict = {
        "id": SCRIPT_ID,
        "file_name": TARGET_FILE,
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": MAX_SCORE,
    }

    try:
        if _PPTX_IMPORT_ERROR is not None:
            result["status"] = "error"
            result["error"] = f"缺少依赖 python-pptx: {_PPTX_IMPORT_ERROR}"
            return result
        base_dir = Path(dir_path)
        if not base_dir.is_dir():
            result["status"] = "error"
            result["error"] = f"目录不存在或不是目录：{dir_path}"
            return result

        target = _locate_target(base_dir)
        if target is None:
            result["status"] = "error"
            result["error"] = f"在目录中未找到 .pptx 文档：{dir_path}"
            return result
        result["file_name"] = target.name

        gate = dimension1_gate(target)
        if not gate.ok or gate.prs is None:
            result["dim1_pass"] = False
            result["dim1_reason"] = gate.detail
            result["total_score"] = 0
            return result

        result["dim1_pass"] = True
        result["dim1_reason"] = ""

        ctx = build_context(target, gate.prs)
        # 加载主题色映射，供 safe_rgb 在遇到 schemeClr 时回退解析
        global _active_theme_map
        _active_theme_map = _load_theme_map(gate.prs)

        total = 0
        max_score = 0
        dim2_items: list[dict] = []
        for checker in CHECKS:
            try:
                cr = checker(ctx)
            except Exception as exc:  # 不让单项异常中断整个评估
                cr = CheckResult(0, getattr(checker, "__name__", "未知检查项"), False, f"检查异常：{exc}")
            max_score += cr.points
            delta = cr.points if cr.passed else 0
            total += delta
            dim2_items.append({
                "rule": cr.name,
                "max_delta": cr.points,
                "delta": delta,
                "hit": cr.passed,
                "detail": "",
            })

        result["dim2_items"] = dim2_items
        result["total_score"] = total
        result["max_score"] = max_score
        return result
    except Exception as exc:
        result["status"] = "error"
        result["error"] = f"{type(exc).__name__}: {exc}"
        return result


if __name__ == "__main__":
    _dir = sys.argv[1] if len(sys.argv) > 1 else str(Path(__file__).resolve().parent)
    print(json.dumps(evaluate(_dir), ensure_ascii=False, indent=2))
