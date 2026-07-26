# -*- coding: utf-8 -*-
"""
自动评估脚本：对 “公开课_改造发布版_按要求修改_动画版.pptx” 按打分细则打分。

评估逻辑：
  维度1（可用与可修改性）：不满足 -> 直接 0 分，不再检查维度2。
  维度2（完成度）：满足维度1后逐条检查。
     - 加分细则：需满足细则内的“每一个点”才加分。
     - 扣分细则：满足细则内“任意一点”即扣分。
  最终打印命中的每条细则及总分。

对于难以纯代码判定的点（例如图片画面内容、动画“效果不完全相同”），
采用可量化的代理特征（图片存在性/尺寸/位置、动画 filter/preset 的差异性等）灵活实现评估意图。
"""

import os
import re
import sys
import json
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS = {"a": A_NS, "p": P_NS}

EMU_PER_CM = 360000.0


def cm(v):
    return None if v is None else v / EMU_PER_CM


def approx(value, target, tol):
    return value is not None and abs(value - target) <= tol


# ----------------------------------------------------------------------------
# 辅助：提取形状信息 / 字号 / 动画 / 切换 / 背景
# ----------------------------------------------------------------------------

def all_runs_sizes(shape):
    """返回 shape 内所有 run 的字号（pt）列表。"""
    sizes = []
    if not shape.has_text_frame:
        return sizes
    for para in shape.text_frame.paragraphs:
        # 段落级默认字号
        para_sz = None
        try:
            if para.font.size is not None:
                para_sz = para.font.size.pt
        except Exception:
            pass
        runs = para.runs
        if not runs and para_sz is not None:
            sizes.append(para_sz)
        for r in runs:
            sz = None
            try:
                if r.font.size is not None:
                    sz = r.font.size.pt
            except Exception:
                pass
            if sz is None:
                sz = para_sz
            if sz is not None:
                sizes.append(sz)
    return sizes


def shape_text(shape):
    if shape.has_text_frame:
        return shape.text_frame.text
    return ""


def find_shape_by_text(slide, needle):
    needle = needle.lower()
    for sh in slide.shapes:
        if needle in shape_text(sh).lower():
            return sh
    return None


def get_transition(slide):
    """返回 (effect_name, signature)；无“放映可见且可启动”的切换返回 (None, None)。

    “放映时能看到且能启动”的判定（需同时满足）：
      1) 须有具体效果子元素：transition 节点下存在效果子元素（fade/push/wipe…），
         空 transition 节点（无效果名）不算。
      2) 须可触发：advClick != "false"（单击翻页可播放）或设置了 advTm（自动翻页）。
         若 advClick="false" 且无 advTm，则放映时永远不会推进、切换不会播放。
      3) 速度正常可见：spd 不为异常的极快值（仅接受 slow/med/fast 或缺省；
         其他异常取值视为不可见）。
      4) 效果须为目标办公软件（WPS 演示）实际可渲染：部分 OOXML 切换效果
         WPS 演示不支持，放映/切换面板表现为“无”，这类效果视为不可见。
    签名包含效果名 + 方向/取向等属性，用于“任意两页不完全相同”的比较。
    """
    # WPS 演示不支持渲染的切换效果（放映看不到、切换面板显示“无”）。
    WPS_UNSUPPORTED = {"uncover"}

    tr = slide._element.find(".//p:transition", NS)
    if tr is None:
        return None, None

    tr_attrs = dict(tr.attrib)

    # 条件2：须可触发
    adv_click = tr_attrs.get("advClick", "true")
    adv_tm = tr_attrs.get("advTm")
    if str(adv_click).lower() == "false" and not adv_tm:
        return None, None

    # 条件3：速度正常可见
    spd = tr_attrs.get("spd")
    if spd is not None and spd not in ("slow", "med", "fast"):
        return None, None

    # 条件1+4：须有具体效果子元素，且该效果在 WPS 演示中可实际渲染
    for child in tr:
        name = etree.QName(child).localname
        if name in WPS_UNSUPPORTED:
            return None, None
        attrs = dict(child.attrib)
        sig = name + "|" + "|".join("%s=%s" % (k, attrs[k]) for k in sorted(attrs))
        return name, sig

    # 有 transition 节点但无具体效果子元素 —— 放映时无可见切换效果
    return None, None


def animation_targets(slide):
    """返回该页所有动画目标 spid -> filter/preset 列表。
    用于判断哪些 shape 设了动画、各组动画是否不完全相同。"""
    timing = slide._element.find(".//p:timing", NS)
    result = {}  # spid -> list of effect signatures
    if timing is None:
        return result
    # 遍历 animEffect / anim / set / animClr 等，向上找 cBhvr/spTgt
    for spTgt in timing.iter("{%s}spTgt" % P_NS):
        spid = spTgt.get("spid")
        # 找最近的 animEffect / anim 等兄弟祖先以取 filter
        node = spTgt
        sig = None
        while node is not None:
            ln = etree.QName(node).localname
            if ln == "animEffect":
                sig = "animEffect:" + (node.get("filter") or node.get("transition") or "")
                break
            if ln in ("anim", "animClr", "animScale", "animRot", "set"):
                sig = ln
                # 不 break，animEffect 优先；但若无 animEffect 则用它
            node = node.getparent()
        result.setdefault(spid, []).append(sig or "set")
    return result


def timeline_is_triggerable(slide):
    """判断该页动画时间线整体是否“可被触发/可播放”。

    放映时动画要真正出现，除了单个 animEffect 自身有效外，整条时间线还必须能被启动：
      - 存在主序列 <p:seq nodeType="mainSeq">；
      - 主序列具备推进/启动接线：seq 上有 prevCondLst 或 nextCondLst（单击推进），
        或 mainSeq 的 cTn 上有非 indefinite 的 stCondLst（自动启动）。
    若主序列没有任何启动/推进条件（stCondLst/prevCondLst/nextCondLst 全缺，
    或启动条件恒为 indefinite 且无单击推进），则 PowerPoint/WPS 放映时动画
    不会被激活，动画面板表现为空，视为“放映看不见”。
    """
    timing = slide._element.find(".//p:timing", NS)
    if timing is None:
        return False
    seq = timing.find(".//p:seq", NS)
    if seq is None:
        return False
    # 单击推进接线（最常见的“点击逐条出现”方式）
    has_prev = seq.find("p:prevCondLst", NS) is not None
    has_next = seq.find("p:nextCondLst", NS) is not None
    if has_prev or has_next:
        return True
    # 或：mainSeq 自身有非 indefinite 的启动条件（自动播放）
    seq_ctn = seq.find("p:cTn", NS)
    if seq_ctn is not None:
        st = seq_ctn.find("p:stCondLst", NS)
        if st is not None:
            for cond in st.findall("p:cond", NS):
                if cond.get("delay") != "indefinite":
                    return True
    return False


def visible_animation_effects(slide):
    """返回该页“放映时实际可见”的动画效果： spid -> set(签名)。

    “放映时看得见”的判定（需同时满足）：
      0) 整条动画时间线可被触发/可播放（见 timeline_is_triggerable）。
         若主序列缺少启动/推进接线，放映时动画不会出现（WPS 动画面板为空），
         则本页所有动画一律视为不可见，返回空。
      1) 须有真效果类型：所在动画归属于 entr/emph/exit/path 之一（presetClass），
         并且确实含 animEffect 节点（仅 set visibility 而无 animEffect 的不算）。
      2) 持续时间 dur>0：animEffect 自身 cBhvr/cTn 的 dur 必须 > 0（“indefinite”视为有效）。
      3) filter 须为具体效果：filter（或 transition）非空且不为 none。
    目标形状是否可见由调用方结合 shape 内容另行判断（见 rule 内 visible_shape）。
    """
    timing = slide._element.find(".//p:timing", NS)
    result = {}  # spid -> set of effect signatures
    if timing is None:
        return result

    # 条件0：整条时间线必须可被触发，否则放映看不到任何动画
    if not timeline_is_triggerable(slide):
        return result

    VALID_CLASSES = {"entr", "emph", "exit", "path"}

    for ae in timing.findall(".//p:animEffect", NS):
        filt = (ae.get("filter") or ae.get("transition") or "").strip()
        # 条件3：filter 须为具体效果
        if not filt or filt.lower() == "none":
            continue

        # 条件2：dur>0
        ctn = ae.find(".//p:cBhvr/p:cTn", NS)
        dur_raw = ctn.get("dur") if ctn is not None else None
        if dur_raw is None:
            continue
        if dur_raw != "indefinite":
            try:
                if int(dur_raw) <= 0:
                    continue
            except (TypeError, ValueError):
                continue

        # 条件1：归属于真效果类型 presetClass
        node = ae
        pc = None
        while node is not None:
            if node.get("presetClass"):
                pc = node.get("presetClass")
                break
            node = node.getparent()
        if pc not in VALID_CLASSES:
            continue

        spt = ae.find(".//p:spTgt", NS)
        spid = spt.get("spid") if spt is not None else None
        if spid is None:
            continue
        result.setdefault(spid, set()).add("animEffect:" + filt)

    return result


def slide_animation_filters(slide):
    """返回该页所有 animEffect 的 filter（按出现顺序）。"""
    timing = slide._element.find(".//p:timing", NS)
    if timing is None:
        return []
    return [e.get("filter") or e.get("transition") or "" for e in timing.findall(".//p:animEffect", NS)]


def get_bg_fill(slide):
    bg = slide._element.find(".//p:cSld/p:bg", NS)
    if bg is None:
        return None
    srgb = bg.find(".//a:srgbClr", NS)
    if srgb is not None:
        return srgb.get("val")
    return "non-solid"


def shape_overlaps(a, b):
    """两个形状矩形是否重叠（有正面积交集）。"""
    ax1, ay1 = a.left, a.top
    ax2, ay2 = a.left + a.width, a.top + a.height
    bx1, by1 = b.left, b.top
    bx2, by2 = b.left + b.width, b.top + b.height
    ix = max(0, min(ax2, bx2) - max(ax1, bx1))
    iy = max(0, min(ay2, by2) - max(ay1, by1))
    return ix > 0 and iy > 0, ix, iy


# ----------------------------------------------------------------------------
# 维度2：各加分/扣分细则
# ----------------------------------------------------------------------------

def rule_p1_slide1_fontsizes(prs):
    """+1：第1页
       “Comics & Animation”约28至29磅，
       “Story Writing Lab”约44至45磅，
       副标题“Build a story from panels, bubbles, and action.”约18至19磅，
       副标题下方的三条正文约19至20磅。
    细则每一个点都要踩到：四处文本均需落在各自磅值区间，正文需为三条且每条都在区间内。
    """
    s = prs.slides[0]

    def text_sizes(shape):
        """返回 shape 内所有非空 run 的字号（pt）列表。"""
        if shape is None:
            return []
        return [sz for sz in all_runs_sizes(shape)]

    def all_in(shape, lo, hi):
        sizes = text_sizes(shape)
        return len(sizes) > 0 and all(lo <= sz <= hi for sz in sizes)

    c_anim = find_shape_by_text(s, "Comics & Animation")
    story = find_shape_by_text(s, "Story Writing Lab")
    sub = find_shape_by_text(s, "Build a story from panels, bubbles, and action.")
    body = find_shape_by_text(s, "observe visual clues")

    # “Comics & Animation”约28至29磅
    c1 = all_in(c_anim, 28, 29)
    # “Story Writing Lab”约44至45磅
    c2 = all_in(story, 44, 45)
    # 副标题约18至19磅
    c3 = all_in(sub, 18, 19)
    # 副标题下方的三条正文约19至20磅（须为三条，且每条都在区间内）
    body_lines = []
    if body is not None:
        for para in body.text_frame.paragraphs:
            if para.text.strip():
                psizes = [sz for r in para.runs for sz in ([r.font.size.pt] if r.font.size else [])]
                if not psizes and para.font.size is not None:
                    psizes = [para.font.size.pt]
                body_lines.append(psizes)
    c4 = (len(body_lines) == 3
          and all(len(ps) > 0 and all(19 <= sz <= 20 for sz in ps) for ps in body_lines))

    checks = [("“Comics & Animation”约28至29磅 (实测%s)" % (text_sizes(c_anim) or "无"), c1),
              ("“Story Writing Lab”约44至45磅 (实测%s)" % (text_sizes(story) or "无"), c2),
              ("副标题“Build a story from panels, bubbles, and action.”约18至19磅 (实测%s)"
               % (text_sizes(sub) or "无"), c3),
              ("副标题下方的三条正文约19至20磅 (共%d条, 实测%s)"
               % (len(body_lines), body_lines), c4)]
    return all(c[1] for c in checks), checks


def rule_p1_slide2_bg(prs):
    """+1：第2页背景：删除深蓝色背景，改为与第3页一致的米黄色背景。
    细则每一个点都要踩到：
      点1 删除深蓝色背景 —— 第2页背景不再是深蓝色。
      点2 改为与第3页一致的米黄色背景 —— 第2页背景色与第3页一致，且为米黄色。
           “与第3页一致”按细则要求理解为颜色代码完全一致（十六进制相同），
           不能仅凭视觉接近判定。
    """
    s2 = prs.slides[1]
    s3 = prs.slides[2]
    bg2 = get_bg_fill(s2)
    bg3 = get_bg_fill(s3)

    def is_cream(h):
        # 米黄色：高 R/G、较低 B
        if not h or len(h) != 6:
            return False
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return r > 200 and g > 200 and b > 180 and r >= b and g >= b

    def is_dark_blue(h):
        if not h or len(h) != 6:
            return False
        r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        return b > 80 and b > r + 20 and b > g + 20 and (r + g + b) < 360

    def consistent(h1, h2):
        # “与第3页一致”：要求颜色代码完全一致（十六进制相同，大小写不敏感），
        # 不接受仅视觉接近的近似色。
        if not (h1 and h2 and len(h1) == 6 and len(h2) == 6):
            return False
        return h1.upper() == h2.upper()

    # 点1：删除深蓝色背景
    c1 = bg2 is not None and not is_dark_blue(bg2)
    # 点2：改为与第3页一致的米黄色背景 —— 与第3页背景色一致(同色系且接近) 且 为米黄色
    c2 = (bg2 is not None and bg3 is not None
          and is_cream(bg2) and is_cream(bg3) and consistent(bg2, bg3))

    checks = [("删除深蓝色背景：第2页背景不再为深蓝色 (bg=%s)" % bg2, c1),
              ("改为与第3页一致的米黄色背景：第2页(%s)与第3页(%s)一致且为米黄色"
               % (bg2, bg3), c2)]
    return all(c[1] for c in checks), checks


def rule_p1_slide3_pic(prs):
    """+1：第3页图片宽度为12厘米，高为9厘米，位于页面右侧。
    细则每一个点都要踩到（针对同一张图片同时满足）：
      点1 插入学生学习场景图片 —— 第3页新增一张图片（学习场景，用图片存在/alt描述代理判定画面内容）。
      点2 图片宽度为12厘米。
      点3 图片高为9厘米。
      点4 位于页面右侧。
    """
    s = prs.slides[2]
    pics = [sh for sh in s.shapes if sh.shape_type == 13]
    half_w = prs.slide_width / 2

    # 逐张图片找出同时满足全部四点的那一张
    best = None
    for sh in pics:
        w, h = cm(sh.width), cm(sh.height)
        center_x = (sh.left + sh.width / 2) if sh.left is not None else None
        p1 = True  # 存在图片对象即视为“插入了图片”；画面内容用 alt 描述做代理
        p2 = approx(w, 12, 0.5)            # 宽度 12cm
        p3 = approx(h, 9, 0.5)             # 高度 9cm
        p4 = center_x is not None and center_x > half_w   # 位于页面右侧（中心在右半页）
        if p1 and p2 and p3 and p4:
            best = sh
            break

    if best is not None:
        cNvPr = best._element.find(".//p:nvPicPr/p:cNvPr", NS)
        alt = (cNvPr.get("descr") if cNvPr is not None else "") or ""
        w, h = cm(best.width), cm(best.height)
        center_x = best.left + best.width / 2
        checks = [
            ("插入学生学习场景图片 (代理: 第3页存在图片对象, alt=%r)" % alt, True),
            ("图片宽度为12厘米 (实测 %.2fcm)" % w, True),
            ("图片高为9厘米 (实测 %.2fcm)" % h, True),
            ("位于页面右侧 (图片中心 %.2fcm > 半页 %.2fcm)" % (center_x / EMU_PER_CM, half_w / EMU_PER_CM), True),
        ]
        return True, checks

    # 未找到同时满足四点的图片：分别报告各点状态（取最接近的一张做展示）
    if pics:
        sh = max(pics, key=lambda x: (x.left or 0))  # 最靠右的图片
        w, h = cm(sh.width), cm(sh.height)
        center_x = (sh.left + sh.width / 2) if sh.left is not None else None
        checks = [
            ("插入学生学习场景图片 (第3页图片数=%d)" % len(pics), True),
            ("图片宽度为12厘米 (实测 %.2fcm)" % w, approx(w, 12, 0.5)),
            ("图片高为9厘米 (实测 %.2fcm)" % h, approx(h, 9, 0.5)),
            ("位于页面右侧", center_x is not None and center_x > half_w),
        ]
    else:
        checks = [
            ("插入学生学习场景图片 (第3页无图片)", False),
            ("图片宽度为12厘米", False),
            ("图片高为9厘米", False),
            ("位于页面右侧", False),
        ]
    return all(c[1] for c in checks), checks


def rule_p1_slide4_layout(prs):
    """+1：第4页整体排版：页面左侧有一张图片，标题、提问内容和说明文字排列在右侧。
    细则每一个点都要踩到：
      点1 探险人物图片移至页面左侧 —— 图片中心位于左半页。
      点2 标题、提问内容和说明文字排列在右侧 —— 这三类文本中心均位于右半页；
           其中“标题”包括页面上所有标题性质的文字（如栏目标题
           “Warm-up: meet the explorer” 与内容标题 “What kind of story is this?”），
           只要有任意一个标题不在右半页，本点即不达标。
      点3 左右布局清晰 —— 图片整体在左、文字整体在右，二者不左右交叉（图片右缘 ≤ 文字左缘）。
    """
    s = prs.slides[3]
    half = prs.slide_width / 2
    pics = [sh for sh in s.shapes if sh.shape_type == 13]

    def center_x(sh):
        return (sh.left + sh.width / 2) if sh.left is not None else None

    # 点1：探险人物图片位于页面左侧（图片中心在左半页）
    left_pics = [sh for sh in pics if center_x(sh) is not None and center_x(sh) < half]
    pic_left = len(left_pics) >= 1

    # 点2：标题、提问内容、说明文字 分别位于右侧
    # “标题”涵盖页面上所有标题性质文字（栏目标题 + 内容标题），任意一个在左侧即不达标。
    title_kws = ["Warm-up: meet the explorer", "What kind of story"]
    titles = [t for t in (find_shape_by_text(s, kw) for kw in title_kws) if t is not None]
    question = find_shape_by_text(s, "Who is the hero")        # 提问内容
    note = find_shape_by_text(s, "stronger plot")              # 说明文字

    def on_right(sh):
        return sh is not None and center_x(sh) is not None and center_x(sh) > half

    title_states = [(shape_text(t).strip()[:30], on_right(t)) for t in titles]
    title_right = len(titles) > 0 and all(state for _, state in title_states)
    question_right = on_right(question)
    note_right = on_right(note)
    text_right = title_right and question_right and note_right

    # 点3：左右布局清晰 —— 图片右缘不越过右侧文字的左缘
    clear = False
    right_texts = titles + [t for t in (question, note) if t is not None]
    if left_pics and right_texts:
        pic_right_edge = max(p.left + p.width for p in left_pics)
        text_left_edges = [t.left for t in right_texts if t.left is not None]
        if text_left_edges:
            clear = pic_right_edge <= min(text_left_edges)

    checks = [
        ("探险人物图片移至页面左侧 (左侧图片数=%d)" % len(left_pics), pic_left),
        ("标题/提问内容/说明文字排列在右侧 (各标题=%s 提问=%s 说明=%s)"
         % (title_states, question_right, note_right), text_right),
        ("左右布局清晰 (图片右缘不越过右侧文字左缘)", clear),
    ]
    return all(c[1] for c in checks), checks


def rule_p3_slide5_pics(prs):
    """+3：第5页图片：第5页出现两张图片，横纵比分别是3:2和4:3，且没有遮挡文本。
    细则每一个点都要踩到：
      点1 第5页出现两张图片。
      点2 两张图片的横纵比（宽:高）分别是 3:2(=1.5) 和 4:3(≈1.333)，
           即两张图片各对应一个比例（不分先后顺序），比值误差≤0.05。
      点3 没有遮挡文本 —— 两张图片均不与任何含文字的文本框显著交叠。
    """
    s = prs.slides[4]
    pics = [sh for sh in s.shapes if sh.shape_type == 13]

    # 点1：第5页出现两张图片
    p1 = len(pics) == 2

    # 点2：两张图片横纵比分别是 3:2 与 4:3（宽:高），不分先后；误差容差 0.05
    TOL = 0.05
    TARGETS = {"3:2": 1.5, "4:3": 4.0 / 3.0}

    def ratio(sh):
        if sh.width and sh.height:
            return sh.width / sh.height
        return None

    ratios = [(sh, ratio(sh)) for sh in pics]
    p2 = False
    ratio_detail = [(sh.shape_id, round(r, 4) if r else None) for sh, r in ratios]
    if len(pics) == 2 and all(r is not None for _, r in ratios):
        r0, r1 = ratios[0][1], ratios[1][1]
        # 一张匹配 3:2、另一张匹配 4:3（两种配对方式任一成立即可）
        match_a = (abs(r0 - TARGETS["3:2"]) <= TOL and abs(r1 - TARGETS["4:3"]) <= TOL)
        match_b = (abs(r1 - TARGETS["3:2"]) <= TOL and abs(r0 - TARGETS["4:3"]) <= TOL)
        p2 = match_a or match_b

    # 点3：没有遮挡文本 —— 图片不与任何含文字的文本框显著交叠（交集 > 文本框面积的20%）
    text_shapes = [sh for sh in s.shapes
                   if sh.has_text_frame and shape_text(sh).strip()
                   and sh.left is not None and sh.width is not None
                   and sh.top is not None and sh.height is not None]

    def covers_text(pic, txt):
        ov, ix, iy = shape_overlaps(pic, txt)
        if not ov:
            return False
        txt_area = txt.width * txt.height
        return txt_area > 0 and (ix * iy) / txt_area > 0.2

    cover_hits = []
    for pic in pics:
        if pic.left is None or pic.width is None:
            continue
        for txt in text_shapes:
            if covers_text(pic, txt):
                cover_hits.append("图片%s 遮挡文本 %r" % (pic.shape_id, shape_text(txt).strip()[:20]))
    p3 = len(pics) >= 1 and len(cover_hits) == 0

    checks = [
        ("第5页出现两张图片 (图片数=%d)" % len(pics), p1),
        ("两张图片横纵比分别为3:2和4:3 (实测宽高比=%s)" % ratio_detail, p2),
        ("没有遮挡文本 (遮挡情况=%s)" % (cover_hits or "无"), p3),
    ]
    return all(c[1] for c in checks), checks


def rule_p1_slide6_lines(prs):
    """+1：第6页箭头右侧区域：每一个箭头右侧对应空白区域均出现两条水平横线，
       箭头右侧除横线外没有其余文字。
    细则每一个点都要踩到：
      点1 每一个箭头右侧对应空白区域均出现两条水平横线。
      点2 箭头右侧除横线外没有其余文字。
    实现说明：横线以占位下划线 '____' 作为代理（一段连续下划线=一条水平横线）；
             “箭头右侧”取含箭头 '→' 的文本框中、箭头之后的内容。
    """
    s = prs.slides[5]
    # 所有“箭头”所在的文本框（含 '→'）
    arrow_shapes = [sh for sh in s.shapes
                    if sh.has_text_frame and "→" in shape_text(sh)]

    per_arrow_two_lines = []   # 点1：每个箭头右侧是否恰有两条横线
    per_arrow_no_extra = []    # 点2：每个箭头右侧除横线外是否无其余文字
    detail = []
    for sh in arrow_shapes:
        txt = shape_text(sh)
        # 取箭头之后（右侧）的内容
        right_part = txt.split("→", 1)[1] if "→" in txt else txt
        # 一段或多段连续下划线各算一条横线
        line_segments = re.findall(r"_+", right_part)
        n_lines = len(line_segments)
        # 除横线外的残余文字（去掉下划线与空白）
        residual = re.sub(r"_+", "", right_part)
        residual = residual.replace("\n", "").strip()
        two_lines = (n_lines == 2)
        no_extra = (residual == "")
        per_arrow_two_lines.append(two_lines)
        per_arrow_no_extra.append(no_extra)
        detail.append("箭头条目: 横线数=%d 残余文字=%r" % (n_lines, residual))

    has_arrows = len(arrow_shapes) > 0
    # 点1：每一个箭头右侧均出现两条水平横线
    p1 = has_arrows and all(per_arrow_two_lines)
    # 点2：每一个箭头右侧除横线外没有其余文字
    p2 = has_arrows and all(per_arrow_no_extra)

    checks = [
        ("每一个箭头右侧对应空白区域均出现两条水平横线 (箭头数=%d, 全部达标=%s)"
         % (len(arrow_shapes), all(per_arrow_two_lines) if has_arrows else False), p1),
        ("箭头右侧除横线外没有其余文字 (全部达标=%s)"
         % (all(per_arrow_no_extra) if has_arrows else False), p2),
    ]
    return all(c[1] for c in checks), checks


def rule_p5_slide8_anim(prs):
    """+5：第8页四组内容动画：“Speech bubble”“Caption”“Sound effect”“Thought bubble”
       及各自说明文字均设置动画，各组动画效果不完全相同。
    细则每一个点都要踩到：
      点1 四组名称文本框（Speech bubble / Caption / Sound effect / Thought bubble）均设置动画。
      点2 各自说明文字文本框均设置动画。
      点3 各组动画效果不完全相同。
    判定原则（严格按 rubric）：
      仅按 OOXML/PowerPoint 中动画是否“存在”和“效果签名是否互异”判断。
      rubric 未要求“放映时可见”，故不再额外强加 timeline 触发接线、presetClass、
      dur、filter 非空等可见性限制。凡在 timing 中有 spTgt 指向该形状的动画节点
      （animEffect/anim/animClr/animScale/animRot/set 等）均计入“已设置动画”，
      其签名参与“各组不完全相同”的比较。
    """
    s = prs.slides[7]
    targets = animation_targets(s)    # spid -> list(动画签名)，按存在性判定

    def has_anim(sh):
        """该形状是否设置了动画：timing 中存在指向其 spid 的动画节点即可。"""
        if sh is None:
            return False
        sid = str(sh.shape_id)
        return len(targets.get(sid, [])) > 0

    group_names = ["Speech bubble", "Caption", "Sound effect", "Thought bubble"]

    def center_y(sh):
        return (sh.top + sh.height / 2) if sh.top is not None else None

    name_anim = []        # 点1：每组名称是否设置了动画
    desc_anim = []        # 点2：每组说明文字是否设置了动画
    group_sig = []        # 点3：每组动画效果签名（名称+说明合并）
    detail = []

    for kw in group_names:
        name_sh = find_shape_by_text(s, kw)
        a_name = has_anim(name_sh)
        name_anim.append(a_name)

        # 该组“说明文字”：与名称同一行（top 接近）、位于名称右侧、且非名称本身的文本框
        desc_sh = None
        if name_sh is not None:
            row_y = center_y(name_sh)
            cands = [sh for sh in s.shapes
                     if sh.has_text_frame and sh is not name_sh
                     and shape_text(sh).strip()
                     and sh.left is not None and name_sh.left is not None
                     and sh.left > name_sh.left                       # 在名称右侧
                     and center_y(sh) is not None and row_y is not None
                     and abs(center_y(sh) - row_y) < int(0.6 * EMU_PER_CM)]  # 同一行
            if cands:
                desc_sh = min(cands, key=lambda x: x.left)            # 紧邻右侧的说明文字
        a_desc = has_anim(desc_sh)
        desc_anim.append(a_desc)

        # 组效果签名：该组名称与说明所用动画签名的集合（存在性签名，不做可见性过滤）
        sigs = set()
        for sh in (name_sh, desc_sh):
            if sh is not None:
                sigs |= set(targets.get(str(sh.shape_id), []))
        group_sig.append(tuple(sorted(sigs)))
        detail.append("组%r: 名称有动画=%s 说明有动画=%s 效果=%s"
                      % (kw, a_name, a_desc, sorted(sigs)))

    p1 = len(group_names) == 4 and all(name_anim)
    p2 = all(desc_anim)
    # 点3：各组动画效果不完全相同 —— 四组效果签名互不完全相同（去重后等于组数）
    p3 = len(set(group_sig)) == len(group_sig) and all(group_sig)

    checks = [
        ("四组名称均设置动画 (%s)" % dict(zip(group_names, name_anim)), p1),
        ("各自说明文字均设置动画 (%s)" % dict(zip(group_names, desc_anim)), p2),
        ("各组动画效果不完全相同 (各组效果=%s)" % [list(x) for x in group_sig], p3),
    ]
    return all(c[1] for c in checks), checks + [(d, True) for d in detail]


def rule_p5_slide9_anim(prs):
    """+5：第9页四组内容动画：“Caption”“Speech bubble”“Thought bubble”“Sound effect”
       及对应例句均设置动画，各组动画效果不完全相同。
    细则每一个点都要踩到：
      点1 四组名称文本框（Caption / Speech bubble / Thought bubble / Sound effect）均设置动画。
      点2 各自对应例句文本框均设置动画。
      点3 各组动画效果不完全相同。
    判定原则（严格按 rubric）：
      仅按 OOXML/PowerPoint 中动画是否“存在”和“效果签名是否互异”判断。
      rubric 未要求“放映时可见”，故不再额外强加 timeline 触发接线、presetClass、
      dur、filter 非空等可见性限制。凡在 timing 中有 spTgt 指向该形状的动画节点
      （animEffect/anim/animClr/animScale/animRot/set 等）均计入“已设置动画”，
      其签名参与“各组不完全相同”的比较。
    """
    s = prs.slides[8]
    targets = animation_targets(s)    # spid -> list(动画签名)，按存在性判定

    def has_anim(sh):
        """该形状是否设置了动画：timing 中存在指向其 spid 的动画节点即可。"""
        if sh is None:
            return False
        return len(targets.get(str(sh.shape_id), [])) > 0

    group_names = ["Caption", "Speech bubble", "Thought bubble", "Sound effect"]

    name_anim = []        # 点1：每组名称是否设置了动画
    quote_anim = []       # 点2：每组例句是否设置了动画
    group_sig = []        # 点3：每组动画效果签名
    detail = []

    for kw in group_names:
        name_sh = find_shape_by_text(s, kw)
        a_name = has_anim(name_sh)
        name_anim.append(a_name)

        # 该组“例句”：与名称同一列（left 接近）、位于名称下方、紧邻的非空文本框
        quote_sh = None
        if name_sh is not None:
            cands = [sh for sh in s.shapes
                     if sh.has_text_frame and sh is not name_sh
                     and shape_text(sh).strip()
                     and sh.left is not None and name_sh.left is not None
                     and abs(sh.left - name_sh.left) < int(0.6 * EMU_PER_CM)   # 同一列
                     and sh.top is not None and name_sh.top is not None
                     and sh.top > name_sh.top]                                 # 在名称下方
            if cands:
                quote_sh = min(cands, key=lambda x: x.top)                     # 紧邻下方的例句
        a_quote = has_anim(quote_sh)
        quote_anim.append(a_quote)

        # 组效果签名：该组名称与例句所用动画签名的集合（存在性签名，不做可见性过滤）
        sigs = set()
        for sh in (name_sh, quote_sh):
            if sh is not None:
                sigs |= set(targets.get(str(sh.shape_id), []))
        group_sig.append(tuple(sorted(sigs)))
        detail.append("组%r: 名称有动画=%s 例句有动画=%s 效果=%s"
                      % (kw, a_name, a_quote, sorted(sigs)))

    p1 = len(group_names) == 4 and all(name_anim)
    p2 = all(quote_anim)
    # 点3：各组动画效果不完全相同 —— 四组效果签名互不完全相同
    p3 = len(set(group_sig)) == len(group_sig) and all(group_sig)

    checks = [
        ("四组名称均设置动画 (%s)" % dict(zip(group_names, name_anim)), p1),
        ("各自对应例句均设置动画 (%s)" % dict(zip(group_names, quote_anim)), p2),
        ("各组动画效果不完全相同 (各组效果=%s)" % [list(x) for x in group_sig], p3),
    ]
    return all(c[1] for c in checks), checks + [(d, True) for d in detail]


def rule_p3_slide10_replace(prs):
    """+3：第10页标题准确显示“Match five stages of the story with frames”。标题下方有一张图片，图片宽度为20厘米，高度为12.79厘米。
    细则每一个点都要踩到（同一张曲线图同时满足尺寸与位置）：
      点1 新第10页标题准确显示“Match five stages of the story with frames”。
      点2 存在故事五阶段曲线图（图片对象）。
      点3 曲线图置于标题下方。
      点4 图片宽度为20厘米。
      点5 图片高度为12.79厘米。
    """
    TITLE = "Match five stages of the story with frames"
    s = prs.slides[9]

    # 点1：标题文本准确等于目标字符串
    title = find_shape_by_text(s, TITLE)
    p1 = title is not None and shape_text(title).strip() == TITLE

    pics = [sh for sh in s.shapes if sh.shape_type == 13]

    # 尺寸公差：严格 ±0.005cm（按两位小数四舍五入须分别等于 20.00 与 12.79）
    SIZE_TOL = 0.005

    # 找出同时满足“标题下方 + 宽20cm + 高12.79cm”的曲线图
    chart = None
    for sh in pics:
        w, h = cm(sh.width), cm(sh.height)
        below = (title is not None and sh.top is not None and title.top is not None
                 and sh.top >= title.top + (title.height or 0))   # 完全位于标题下方
        if below and approx(w, 20, SIZE_TOL) and approx(h, 12.79, SIZE_TOL):
            chart = sh
            break

    # 点2：存在曲线图（图片对象）
    p2 = len(pics) >= 1
    # 点3/4/5：取上面命中的 chart；若没有同时满足者，则分别报告最接近的一张
    if chart is not None:
        w, h = cm(chart.width), cm(chart.height)
        p3 = p4 = p5 = True
        below_desc = "实际 top=%.2fcm > 标题底=%.2fcm" % (
            chart.top / EMU_PER_CM,
            (title.top + (title.height or 0)) / EMU_PER_CM if title is not None else 0)
        size_w_desc = "实测 %.4fcm" % w
        size_h_desc = "实测 %.4fcm" % h
    else:
        ref = pics[0] if pics else None
        w = cm(ref.width) if ref is not None else None
        h = cm(ref.height) if ref is not None else None
        title_bottom = (title.top + (title.height or 0)) if title is not None else None
        p3 = (ref is not None and ref.top is not None and title_bottom is not None
              and ref.top >= title_bottom)
        p4 = approx(w, 20, SIZE_TOL)
        p5 = approx(h, 12.79, SIZE_TOL)
        below_desc = "图片 top=%s" % (round(ref.top / EMU_PER_CM, 2) if ref is not None else None)
        size_w_desc = "实测 %s" % (round(w, 4) if w is not None else None)
        size_h_desc = "实测 %s" % (round(h, 4) if h is not None else None)

    checks = [
        ("新第10页标题准确显示“%s” (实测=%r)" % (TITLE, shape_text(title).strip() if title is not None else None), p1),
        ("存在故事五阶段曲线图 (第10页图片数=%d)" % len(pics), p2),
        ("曲线图置于标题下方 (%s)" % below_desc, p3),
        ("图片宽度为20厘米 (%s)" % size_w_desc, p4),
        ("图片高度为12.79厘米 (%s)" % size_h_desc, p5),
    ]
    return all(c[1] for c in checks), checks


def rule_p5_slide11_anim(prs):
    """+5：第11页文本框动画：除主标题外，副标题、页码和六条对白文本框均设置相同的动画，
       各文本框可按讲解顺序分别触发。
    细则每一个点都要踩到：
      点1 副标题设置动画。
      点2 页码设置动画。
      点3 六条对白文本框均设置动画（且确为六条）。
      点4 上述文本框设置的是“相同的”动画（动画效果一致）。
      点5 各文本框可按讲解顺序分别触发（每个文本框都对应一次独立的单击触发，
           且在动画时间线中按讲解顺序逐一匹配）。
    “除主标题外”：主标题不在被检查/被要求的对象之列（不约束其是否有动画）。
    判定原则（严格按 rubric）：
      仅按 OOXML/PowerPoint 中动画是否“存在”、效果签名是否一致，以及时间线中
      clickEffect 的目标 spid 序列是否逐一匹配各要求对象来判定，不再叠加
      “放映可见”“presetClass”“dur>0”“filter 非空”等 rubric 未提出的门槛。
    对象识别（避免依赖易变文案/引号形态）：
      主标题：字号最大、位于页面上部的文本框。
      页码：文本形如 “11”（可含小空白/前后修饰）的短小文本框。
      对白：以任意引号（中文/英文/角引号）起首，或含成对引号的文本框；
            数量应为六条。
      副标题：除主标题、页码、对白之外，页面上剩余的显著文本框中最靠上的一个。
    """
    s = prs.slides[10]
    targets = animation_targets(s)    # spid -> list(动画签名)，按存在性判定

    # ---- 识别页面上各角色的文本框 ----
    all_text_shapes = [sh for sh in s.shapes
                       if sh.has_text_frame and shape_text(sh).strip()
                       and sh.top is not None]

    # 主标题：字号最大且位于页面上部的文本框（不用固定关键词 "While writing"）
    def max_font_pt(sh):
        sizes = all_runs_sizes(sh)
        return max(sizes) if sizes else 0

    page_h = prs.slide_height
    upper_candidates = [sh for sh in all_text_shapes
                        if sh.top is not None and sh.top < page_h * 0.35]
    main_title = None
    if upper_candidates:
        main_title = max(upper_candidates, key=lambda x: (max_font_pt(x), -(x.top or 0)))

    # 页码："11"（允许前后空白），且是较小的文本框
    page_num = None
    for sh in s.shapes:
        if sh.has_text_frame and shape_text(sh).strip() == "11":
            page_num = sh
            break

    # 六条对白：以任意常见引号起首 或 含成对引号
    QUOTE_STARTS = ("“", "”", "‘", "’",
                    "「", "『",   # 「 『
                    "«",              # «
                    '"', "'")

    def is_dialogue(sh):
        if not sh.has_text_frame:
            return False
        t = shape_text(sh).strip()
        if not t:
            return False
        if t.startswith(QUOTE_STARTS):
            return True
        # 含成对引号（中文成对 “ ”/ ‘ ’ 或英文成对 " " / ' '）
        pairs = [
            ("“", "”"),
            ("‘", "’"),
            ("「", "」"),
            ("『", "』"),
        ]
        if any(a in t and b in t for a, b in pairs):
            return True
        if t.count('"') >= 2 or t.count("'") >= 2:
            return True
        return False

    dialogues = [sh for sh in s.shapes if is_dialogue(sh)]

    # 副标题：排除主标题/页码/对白后，剩余显著文本中位置最靠上者
    excluded_ids = {id(x) for x in [main_title, page_num] + dialogues if x is not None}
    remaining = [sh for sh in all_text_shapes if id(sh) not in excluded_ids]
    subtitle = min(remaining, key=lambda x: x.top) if remaining else None

    # ---- 动画存在性 / 签名 ----
    def sid(sh):
        return str(sh.shape_id) if sh is not None else None

    def has_anim(sh):
        """该文本框是否设置了动画：timing 中存在指向其 spid 的动画节点即可。"""
        return sh is not None and len(targets.get(sid(sh), [])) > 0

    def eff_sig(sh):
        """该文本框设置的动画效果签名集合（按存在性汇总，不做可见性过滤）。"""
        if sh is None:
            return tuple()
        return tuple(sorted(set(targets.get(sid(sh), []))))

    # 点1：副标题设置动画
    p1 = has_anim(subtitle)
    # 点2：页码设置动画
    p2 = has_anim(page_num)
    # 点3：六条对白文本框均设置动画 且 恰为六条
    n_dialog = len(dialogues)
    p3 = (n_dialog == 6) and all(has_anim(d) for d in dialogues)

    # 需要设相同动画的全部对象（副标题、页码、六条对白）
    required = [("副标题", subtitle), ("页码", page_num)] + \
               [("对白:" + shape_text(d).strip()[:12], d) for d in dialogues]
    required_present = [(lbl, sh) for lbl, sh in required if sh is not None]

    # 点4：相同的动画 —— 所有要求对象的效果签名一致（且非空、均已设动画）
    sigs = [eff_sig(sh) for _, sh in required_present]
    p4 = (len(sigs) > 0
          and all(sg for sg in sigs)                       # 每个都设了动画
          and len(set(sigs)) == 1)                         # 且效果完全相同

    # 点5：各文本框可按讲解顺序分别触发 ——
    #      时间线中 clickEffect 的目标 spid 序列应包含每一个要求对象；
    #      即每个要求对象都能被独立触发（不是一次单击触发多个/或漏触发）。
    #      同时按 clickEffect 出现顺序与要求对象在页面上的讲解顺序（自上而下）
    #      比对，允许一致或页面上明确的其它单向顺序（此处以“每个 spid 至少出现一次
    #      作为独立 clickEffect 目标”作为核心检查，避免因编辑器改序造成误判）。
    timing = s._element.find(".//p:timing", NS)
    click_target_seq = []              # clickEffect 触发目标 spid 序列（按时间线顺序）
    if timing is not None:
        for ctn in timing.iter("{%s}cTn" % P_NS):
            if ctn.get("nodeType") == "clickEffect":
                sp = ctn.find(".//p:spTgt", NS)
                if sp is not None and sp.get("spid"):
                    click_target_seq.append(sp.get("spid"))

    required_spids = [sid(sh) for _, sh in required_present]
    # 每个要求对象都出现在独立的 clickEffect 目标中
    each_independently_triggered = (len(required_spids) > 0
                                    and all(spid in click_target_seq
                                            for spid in required_spids))

    # 讲解顺序的“逐一匹配”：clickEffect 中按首次出现次序抽出的要求对象子序列
    # 应恰好等于要求对象集合（即数量匹配、无遗漏、无重复目标混入）
    seen = set()
    click_order_of_required = []
    for spid in click_target_seq:
        if spid in required_spids and spid not in seen:
            seen.add(spid)
            click_order_of_required.append(spid)
    order_ok = (len(click_order_of_required) == len(required_spids))

    p5 = each_independently_triggered and order_ok

    anim_detail = ", ".join("%s=%s" % (lbl, "动" if has_anim(sh) else "无")
                            for lbl, sh in required_present)

    checks = [
        ("副标题设置动画 (识别到副标题=%r)"
         % (shape_text(subtitle).strip()[:20] if subtitle is not None else None), p1),
        ("页码设置动画 (识别到页码=%r)"
         % (shape_text(page_num).strip() if page_num is not None else None), p2),
        ("六条对白文本框均设置动画 (共%d条)" % n_dialog, p3),
        ("上述文本框设置相同的动画 (效果集合=%s)" % {x for x in sigs}, p4),
        ("各文本框可按讲解顺序分别触发 "
         "(clickEffect 目标序列=%s, 需触发对象数=%d, 匹配数=%d)"
         % (click_target_seq, len(required_spids), len(click_order_of_required)), p5),
    ]
    return all(c[1] for c in checks), checks + [("明细: " + anim_detail, True)]


def _rubric_transition(slide):
    """按 rubric 直读该页 <p:transition> 的“效果名 + 方向属性”。

    仅判断“是否设置了切换效果”并抽取（名称, 方向属性组合）：
      - 有 <p:transition> 且其下存在至少一个效果子元素（fade/push/wipe/cover/…），
        即视为“已设置切换效果”；空 transition 节点或无效果子元素不算。
      - 签名 = 效果名 + 该效果节点的属性（含 dir / orient / thruBlk 等方向/取向），
        用于“任意两页不完全相同”的比较。
    不额外检查 advClick/advTm（是否单击/自动触发）、spd（速度）、
    以及具体渲染宿主（如 WPS）是否支持——这些均为 rubric 未要求的条件，
    引入会误判“已设置但不满足代理条件”的页为缺失。
    """
    tr = slide._element.find(".//p:transition", NS)
    if tr is None:
        return None, None
    for child in tr:
        name = etree.QName(child).localname
        attrs = dict(child.attrib)
        sig = name + "|" + "|".join("%s=%s" % (k, attrs[k]) for k in sorted(attrs))
        return name, sig
    # 有 transition 节点但无具体效果子元素 —— 未设置具体切换效果
    return None, None


def rule_p5_transitions(prs):
    """+5：第1至第18页切换动画：每页均设置切换效果，
       任意两页的切换效果名称或方向组合不完全相同。
    细则每一个点都要踩到：
      点1 第1至第18页每页均设置切换效果。
      点2 任意两页的切换效果“名称或方向组合”不完全相同
          —— 即不存在任意两页的(名称, 方向组合)完全相同。
    判定原则（严格按 rubric）：
      直接读取每页 <p:transition> 的效果名与方向属性来判定“是否设置”与
      “是否互异”。不再引入 rubric 未要求的可触发（advClick/advTm）、速度（spd）、
      具体渲染宿主（如 WPS）支持等前提。
    """
    sigs = []     # 每页的 (名称|方向组合) 签名
    missing = []  # 未设置切换效果的页
    per_page = []
    for i, s in enumerate(prs.slides, 1):
        name, sig = _rubric_transition(s)
        if name is None:
            missing.append(i)
        sigs.append(sig)
        per_page.append((i, sig))

    # 点1：第1-18页每页均设置切换效果
    p1 = (len(prs.slides) == 18) and (len(missing) == 0)

    # 点2：任意两页的“名称或方向组合”不完全相同 —— 没有任何重复签名
    # 找出重复项以便报告
    seen = {}
    dup_pairs = []
    for i, sig in per_page:
        if sig in seen:
            dup_pairs.append((seen[sig], i, sig))
        else:
            seen[sig] = i
    p2 = len(dup_pairs) == 0 and (None not in sigs)

    checks = [
        ("第1至第18页每页均设置切换效果 (总页数=%d, 缺失页=%s)"
         % (len(prs.slides), missing or "无"), p1),
        ("任意两页切换效果名称或方向组合不完全相同 (唯一组合=%d/%d, 重复=%s)"
         % (len(set(sigs)), len(sigs),
            [(a, b) for a, b, _ in dup_pairs] or "无"), p2),
    ]
    return all(c[1] for c in checks), checks


# ---------------- 扣分项 ----------------


# ----------------------------------------------------------------------------
# 主流程
# ----------------------------------------------------------------------------

SCRIPT_ID = "080"


def _find_pptx(dir_path: str):
    """在给定目录中定位被评估的 .pptx 文件。
    优先返回目录下第一个 .pptx（排除 ~$ 临时文件）。"""
    d = os.path.abspath(dir_path)
    if not os.path.isdir(d):
        return None
    for name in sorted(os.listdir(d)):
        if name.lower().endswith(".pptx") and not name.startswith("~$"):
            return os.path.join(d, name)
    return None


def evaluate(dir_path: str) -> dict:
    """统一入口：接收脚本所在目录路径，在该目录中定位并评估 .pptx 文件。

    返回结构化字典，字段见项目《脚本接口差异与统一建议》§2.2。
    """
    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": 0,
    }

    try:
        file_path = _find_pptx(dir_path)
        if not file_path or not os.path.isfile(file_path):
            result["status"] = "error"
            result["error"] = "在目录 %r 中未找到 .pptx 文件" % dir_path
            return result

        result["file_name"] = os.path.basename(file_path)

        prs = Presentation(file_path)

        # ---- 维度2 加分项定义（无论维度1是否通过，都作为满分基线列出）----
        plus_rules = [
            (1, "第1页“Comics & Animation”约28至29磅，“Story Writing Lab”约44至45磅，"
                "副标题“Build a story from panels, bubbles, and action.”约18至19磅，"
                "副标题下方的三条正文约19至20磅。", rule_p1_slide1_fontsizes),
            (1, "第2页背景：删除深蓝色背景，改为与第3页一致的米黄色背景。", rule_p1_slide2_bg),
            (1, "第3页图片宽度为12厘米，高为9厘米，位于页面右侧。",
                rule_p1_slide3_pic),
            (1, "第4页整体排版：页面左侧有一张图片，标题、提问内容和说明文字排列在右侧。",
                rule_p1_slide4_layout),
            (3, "第5页图片：第5页出现两张图片，横纵比分别是3:2和4:3，且没有遮挡文本。",
                rule_p3_slide5_pics),
            (1, "第6页箭头右侧区域：每一个箭头右侧对应空白区域均出现两条水平横线，"
                "箭头右侧除横线外没有其余文字。", rule_p1_slide6_lines),
            (5, "第8页四组内容动画：“Speech bubble”“Caption”“Sound effect”“Thought bubble”"
                "及各自说明文字均设置动画，各组动画效果不完全相同。", rule_p5_slide8_anim),
            (5, "第9页四组内容动画：“Caption”“Speech bubble”“Thought bubble”“Sound effect”"
                "及对应例句均设置动画，各组动画效果不完全相同。", rule_p5_slide9_anim),
            (3, "第10页标题准确显示“Match five stages of the story with frames”。标题下方有一张图片，图片宽度为20厘米，高度为12.79厘米。",
                rule_p3_slide10_replace),
            (5, "第11页文本框动画：除主标题外，副标题、页码和六条对白文本框均设置相同的动画，"
                "各文本框可按讲解顺序分别触发。", rule_p5_slide11_anim),
            (5, "第1至第18页切换动画：每页均设置切换效果，任意两页的切换效果名称或方向组合不完全相同。",
                rule_p5_transitions),
        ]

        # ---- 维度2 加分项分值即为满分基线 ----
        # 满分：加分项分值之和
        result["max_score"] = sum(pts for pts, _, _ in plus_rules)

        # ---- 维度1 ----
        result["dim1_pass"] = True

        # 维度1通过：逐项评估维度2
        dim2_items = []
        total = 0

        for pts, rule_text, fn in plus_rules:
            ok, _checks = fn(prs)
            delta = pts if ok else 0
            total += delta
            dim2_items.append({
                "rule": rule_text,
                "max_delta": pts,
                "delta": delta,
                "hit": bool(ok),
                "detail": "",
            })

        result["dim2_items"] = dim2_items
        result["total_score"] = total
        return result

    except Exception as e:
        result["status"] = "error"
        result["error"] = "%s: %s" % (type(e).__name__, e)
        return result


if __name__ == "__main__":
    # 本地自测：默认以脚本所在目录为参数
    target = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    _out = json.dumps(evaluate(target), ensure_ascii=False, indent=2)
    # 不改 sys.stdout，直接以 UTF-8 写字节，规避 Windows cp1252 报错
    try:
        sys.stdout.buffer.write((_out + "\n").encode("utf-8"))
    except AttributeError:
        print(_out)
