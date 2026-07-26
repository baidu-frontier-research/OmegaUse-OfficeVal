# -*- coding: utf-8 -*-
"""
自动评估 "第二课时_蒙古语版.pptx"

评估逻辑：
  维度1（可用与可修改性）—— 一票否决。任何一条不满足 -> 直接 0 分，不再检查维度2。
  维度2（完成度）—— 在维度1通过后逐条检查：
      得分点：13页中每一页正文为蒙古语，每满足一页 +3。
      扣分点：命中任意一条扣分项即累加该项分值（负分）。
  最后打印：命中的点 + 最终得分。

实现说明（"灵活变通"之处）：
  - 字号：源文件被整体缩放保存，存储字号 ≈ 细则字号 × 0.86 或 × 0.92。
    因此字号判定为：存储字号 / 细则字号 是否落在 [0.84, 0.94] 区间附近（含原始 1.0）。
    只要落在该区间即视为"字号正确"。
  - 颜色：黑=#22313F、灰=#6B7C84/#22313F系深灰、深绿=#0B5965（深绿色）。
    判定按色系归类（black / gray / darkgreen / green），与细则要求的色系比较。
  - 文字遮挡/显示不全：通过 normAutofit 的 fontScale 是否 <100% 来近似判断
    （PowerPoint 在文本溢出文本框时会自动缩小，fontScale<100 即说明放不下）。
  - 文本定位：按"去空白后包含关系"匹配细则中引用的蒙古语片段，避免空格/换行差异。
"""

import os
import re
import sys

from pptx import Presentation

SCRIPT_ID = "079"

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


# --------------------------------------------------------------------------
# 工具函数
# --------------------------------------------------------------------------
def norm(s):
    """去掉所有空白字符，便于片段匹配。"""
    if not s:
        return ""
    return re.sub(r"\s+", "", s)


def is_chemical_or_symbol(text):
    """判断一段文本是否为化学式/化学符号/纯符号，而非自然语言正文。
    例如 "—CH3""CH2CH3""—CnH2n+1""C""—R" 等。
    判据：去空白后只由 化学元素相关拉丁字母(C/H/R/n等) + 数字 + 化学符号(—-+()) 组成。"""
    s = re.sub(r"\s+", "", text or "")
    if not s:
        return True
    # 只允许化学式常见字符；出现任何其它字母(尤其西里尔)即视为正文
    return bool(s) and re.fullmatch(r"[CHRnOcho0-9—\-+()（）·,，、]+", s) is not None


def is_mongolian(text):
    """判断一段文字主体是否为蒙古语（西里尔蒙文）—— 宽松版。
    仅用于图片 OCR 兜底、遮挡检测等"识别是否为蒙古语文本框"的场景，
    对正文语言判定请使用 is_mongolian_strict。"""
    if not text or not text.strip():
        return False
    cyr = sum(1 for ch in text if "Ѐ" <= ch <= "ӿ")
    letters = sum(1 for ch in text if ch.isalpha())
    if letters == 0:
        return False
    return cyr / letters >= 0.5


# 化学式/化学符号常用的拉丁字母。判定"正文皆为蒙古语"时把它们剔除，
# 避免"C атом"、"H атом"等把 C/H 计入非蒙古语字母而误判。
_CHEM_LATIN = set("CHRnOcho")


def is_mongolian_strict(text: str, min_cyr: int = 2, ratio_min: float = 0.9) -> bool:
    """严格判定：一段正文是否"皆为蒙古语"（西里尔蒙文）。
    规则：
      1. 出现 CJK / 谚文 / 阿拉伯 / 希腊 等其它书写系统字母 —— 直接否；
      2. 剔除化学元素类拉丁字母(C/H/R/n/O 等)后，比较西里尔字母与其余拉丁字母的占比，
         西里尔占比必须 >= ratio_min（默认 0.9）；
      3. 西里尔字母总数必须 >= min_cyr，防止极短文本仅含单个西里尔字符即通过。
    """
    if not text or not text.strip():
        return False
    for ch in text:
        if not ch.isalpha():
            continue
        code = ord(ch)
        # 允许：西里尔基本区 (Ѐ-ӿ, U+0400-U+04FF) 与补充区 (U+0500-U+052F)
        if 0x0400 <= code <= 0x052F:
            continue
        # 允许：基本拉丁字母（可能是化学元素/符号），保留到下方比例判断阶段
        if 0x0041 <= code <= 0x005A or 0x0061 <= code <= 0x007A:
            continue
        # 其它书写系统字母（CJK / 希腊 / 阿拉伯 / 韩文 等）—— 非蒙古语正文
        return False
    cyr = 0
    other_letters = 0
    for ch in text:
        if not ch.isalpha():
            continue
        if "Ѐ" <= ch <= "ӿ":
            cyr += 1
        elif ch in _CHEM_LATIN:
            # 化学元素类拉丁字母：不计入语言判定的分子分母
            continue
        else:
            other_letters += 1
    if cyr < min_cyr:
        return False
    denom = cyr + other_letters
    if denom == 0:
        return False
    return (cyr / denom) >= ratio_min


def classify_color(rgb):
    """把 RGB 颜色归到色系：black / gray / darkgreen / green / other。"""
    if rgb is None:
        return None
    r, g, b = (rgb >> 16) & 0xFF, (rgb >> 8) & 0xFF, rgb & 0xFF
    mx, mn = max(r, g, b), min(r, g, b)
    chroma = mx - mn                 # 饱和度（绝对色度）
    light = (r + g + b) / 3.0

    # 低色度 -> 无明显色相，按明度分黑/灰
    if chroma <= 45:
        if light <= 75:
            return "black"           # 极深，如 22313F（深灰蓝，作黑色正文）
        if light <= 180:
            return "gray"            # 中等明度灰，如 6B7C84
        return "other"

    # 有明显色相：判断是否绿/青系（g、b 明显高于 r，即青绿/墨绿）
    teal_green = (g - r) >= 25 and (b - r) >= 20 and r <= 80
    if teal_green:
        return "darkgreen" if light < 140 else "green"   # 如 0B5965 偏暗 -> 深绿
    return "other"


SIZE_TOL_PT = 0.05  # 字号绝对容差（磅）。仅允许极小浮点/取整误差，
                    # 办公软件实际显示的字号必须等于细则字号，不接受 0.86 / 0.92 等缩放。


def size_matches(stored_pt, spec_pt):
    """字号是否严格等于细则字号（仅允许 ±0.05pt 的浮点/取整误差）。"""
    if stored_pt is None:
        return False
    return abs(stored_pt - spec_pt) <= SIZE_TOL_PT


# --------------------------------------------------------------------------
# 提取每页形状的文本与样式
# --------------------------------------------------------------------------
class RunInfo:
    def __init__(self, text, size_pt, font, color_cls):
        self.text = text
        self.size_pt = size_pt
        self.font = font
        self.color_cls = color_cls


def _endpara_style(para):
    """当 run 的 rPr 为空时，回退到该段 endParaRPr 上携带的样式。
    返回 (size_pt, font, color_cls)。"""
    ep = para._p.find("{%s}endParaRPr" % A)
    if ep is None:
        return None, None, None
    sz = ep.get("sz")
    size_pt = (int(sz) / 100.0) if sz else None
    latin = ep.find("{%s}latin" % A)
    font = latin.get("typeface") if latin is not None else None
    color_cls = None
    srgb = ep.find("{%s}solidFill/{%s}srgbClr" % (A, A))
    if srgb is not None and srgb.get("val"):
        try:
            color_cls = classify_color(int(srgb.get("val"), 16))
        except Exception:
            color_cls = None
    return size_pt, font, color_cls


def shape_runs(shape):
    """返回该形状内所有非空 run 的样式信息。"""
    out = []
    if not shape.has_text_frame:
        return out
    for para in shape.text_frame.paragraphs:
        ep_size, ep_font, ep_color = _endpara_style(para)
        for run in para.runs:
            if not run.text.strip():
                continue
            f = run.font
            size_pt = f.size.pt if f.size is not None else None
            color_cls = None
            try:
                if f.color is not None and f.color.type is not None:
                    # RGBColor 是 str 子类（六位十六进制），需按 16 进制转 int
                    color_cls = classify_color(int(str(f.color.rgb), 16))
            except Exception:
                color_cls = None
            font = f.name
            # run 自身未显式声明样式时，回退到段落 endParaRPr（PowerPoint 实际渲染所用）
            if size_pt is None:
                size_pt = ep_size
            if font is None:
                font = ep_font
            if color_cls is None:
                color_cls = ep_color
            out.append(RunInfo(run.text, size_pt, font, color_cls))
    return out


def shape_has_autofit(shape):
    """文本框是否开启 normAutofit（自动调整字号）。
    开启后办公软件实际显示的字号会偏离存储字号（本文件即被此手法放大显示），
    因此凡目标文本框开了 normAutofit，一律视为"字号不可信/不达标"。"""
    if not shape.has_text_frame:
        return False
    return shape.text_frame._txBody.find(
        "{%s}bodyPr/{%s}normAutofit" % (A, A)) is not None


def shape_text_overflow(shape):
    """文字显示不全：文本框开启了 normAutofit 缩放（fontScale<100%，PowerPoint 因放不下而缩字），
    或文本形状超出幻灯片边界被裁切。返回 True 表示该形状文字显示不全。"""
    if not shape.has_text_frame:
        return False
    body = shape.text_frame._txBody
    naf = body.find("{%s}bodyPr/{%s}normAutofit" % (A, A))
    if naf is not None:
        fs = naf.get("fontScale")
        if fs is not None:
            try:
                if (int(fs) / 1000.0) < 99.0:  # fontScale 以千分之一存储
                    return True
            except Exception:
                pass
    return False


# --------------------------------------------------------------------------
# 维度2 扣分项规则定义
#   每条： (页码, [蒙古语片段...], 要求色系, 要求字号pt, 扣分值, 描述)
#   命中条件：在该页找到"包含任一片段"的形状，且其文字 颜色/字号 不满足要求。
# --------------------------------------------------------------------------
FONT_REQ = "微软雅黑"

RULES = [
    # 第一页
    (1, ["“Нүүрсустөрөгчийн бүлэг”-ээс эхэлж, энгийн нэршил, системчилсэн нэршил болон изомер бичих аргыг алхам алхмаар эзэмшинэ."],
     "black", 16, -1,
     "第一页“Нүүрсустөрөгчийн бүлэг”-ээс эхэлж…эзэмшинэ.文字不是黑色、微软雅黑、16磅"),
    (1, ["Сурах маршрут: ойлголтын суурь → дүрэм задлах → бичих дадлага → дүгнэн бататгах"],
     "gray", 14, -1,
     "第一页“Сурах маршрут: ойлголтын суурь → дүрэм задлах → бичих дадлага → дүгнэн бататгах”文字不是灰色、微软雅黑、14磅"),
    (1, ["Холбоос: алканы нэршлийн сэдэвт орно　Дараа нь: орлуулагч болж чадах хэсгийг танина"],
     "gray", 9.5, -1,
     "第一页“Холбоос: алканы нэршлийн сэдэвт орно　Дараа нь: орлуулагч болж чадах хэсгийг танина”文字不是灰色、微软雅黑、9.5磅"),
    # 第二页
    (2, ["Энэ хуудсанд ойлголтын орцыг үргэлжлүүлж, илүү товч хэлбэрээр тайлбарлав."],
     "gray", 12.5, -1,
     "第二页“Энэ хуудсанд ойлголтын орцыг үргэлжлүүлж, илүү товч хэлбэрээр тайлбарлав.”文字不是灰色、微软雅黑、12.5磅"),
    (2, ["Нүүрсустөрөгчийн молекулаас нэг H атом салгахад үлдэх саармаг хэсгийг —R гэж бичнэ. Зураас нь ковалент холбоо үүсгэх нэг байр байгааг илэрхийлнэ.",
         "Ханасан гинжин алкил бүлгийг —CnH2n+1 гэж илэрхийлнэ. Жишээ нь метил, этил, пропил."],
     "black", 14, -1,
     "第二页“Нүүрсустөрөгчийн молекулаас…илэрхийлнэ.”及“Ханасан гинжин алкил бүлгийг —CnH2n+1 гэж…пропил.”文字不是黑色、微软雅黑、14磅"),
    # 第三页
    (3, ["Энгийн нэршил: эхлээд C-г тоолж, дараа нь “алкан” нэрлэнэ"],
     "darkgreen", 26, -1,
     "第三页“Энгийн нэршил: эхлээд C-г тоолж, дараа нь ”алкан“ нэрлэнэ”文字不是深绿色、微软雅黑、26磅"),
    (3, ["Энэ хуудсанд “хэсэг таних”-ыг үргэлжлүүлж, бүтэн алканыг нүүрстөрөгчийн тоогоор үндсэн нэрлэнэ"],
     "gray", 12.5, -1,
     "第三页“Энэ хуудсанд ”хэсэг таних“-ыг үргэлжлүүлж, бүтэн алканыг нүүрстөрөгчийн тоогоор үндсэн нэрлэнэ”文字不是灰色、微软雅黑、12.5磅"),
    (3, ["Метан, этан, пропан, бутан… гэсэн дарааллаар C атомын тоог илэрхийлж алканы нэр өгнө.",
         "Нүүрстөрөгчийн тоог шууд тоогоор илэрхийлнэ, жишээ нь “додекан”",
         "Молекулын томьёо ижил боловч C араг яс өөр бол n-, изо-, нео- зэрэг бүтцийг ялгана."],
     "black", 14, -1,
     "第三页“Метан, этан, пропан, бутан…өгнө.”、“Нүүрстөрөгчийн тоог шууд тоогоор илэрхийлнэ, жишээ нь ”додекан“、”Молекулын томьёо ижил боловч C араг яс өөр бол n-, изо-, нео- зэрэг бүтцийг ялгана.“文字不是黑色、微软雅黑、14磅"),
    # 第四页
    (4, ["Энэ хуудас энгийн нэршлийг үргэлжлүүлж, C5H12-ийн гурван түгээмэл C араг ясыг үзүүлнэ"],
     "gray", 11.5, -1,
     "第四页“Энэ хуудас энгийн нэршлийг үргэлжлүүлж, C5H12-ийн гурван түгээмэл C араг ясыг үзүүлнэ”文字不是灰色、微软雅黑、11.5磅"),
    (4, ["Бүх C атом тасралтгүй нэг гинж үүсгэнэ. Нэр нь n-пентан.",
         "Гол гинж дээр нэг метил гарна. Нэрийг изопентан гэж бичиж болно.",
         "Төв C нь дөрвөн C бүлэгтэй холбогдоно. Нэр нь неопентан."],
     "black", 13.5, -1,
     "第四页“Бүх C атом тасралтгүй нэг гинж үүсгэнэ. Нэр нь n-пентан.”、“Гол гинж дээр нэг метил гарна. Нэрийг изопентан гэж бичиж болно.”、“Төв C нь дөрвөн C бүлэгтэй холбогдоно. Нэр нь неопентан.”文字不是黑色、微软雅黑、13.5磅"),
    # 第五页
    (5, ["Өмнөх хуудсанд C араг ясны ялгааг харлаа. Энэ хуудсанд төвөгтэй бүтцийг системийн дүрмээр авч үзнэ"],
     "gray", 12.5, -1,
     "第五页“Өмнөх хуудсанд C араг ясны ялгааг харлаа. Энэ хуудсанд төвөгтэй бүтцийг системийн дүрмээр авч үзнэ”文字不是灰色、微软雅黑、12.5磅"),
    (5, ["C атом хамгийн олонтой тасралтгүй гинжийг гол гинжээр түрүүлж сонгоно.",
         "Хамгийн урт гинж нэгээс олон бол орлуулагч илүү олон холбогдсон гинжийг сонгоно.",
         "Гол гинжийн нэрийг түүний C атомын тоогоор тогтооно; салаа гинжийг орлуулагч болгон гол нэрийн өмнө бичнэ."],
     "black", 17, -1,
     "第五页“C атом хамгийн олонтой тасралтгүй гинжийг гол гинжээр түрүүлж сонгоно.”、“Хамгийн урт гинж нэгээс олон бол орлуулагч илүү олон холбогдсон гинжийг сонгоно.”、“Гол гинжийн нэрийг түүний C атомын тоогоор тогтооно; салаа гинжийг орлуулагч болгон гол нэрийн өмнө бичнэ.”文字不是黑色、微软雅黑、17磅"),
    # 第六页
    (6, ["Орлуулагчид ойр талаас дугаарлаж, салаа гинжид бага байрлалын дугаар өгнө. Дугаарлалт нь дурын биш; нэрийг аль болох тодорхой, цорын ганц болгох зорилготой."],
     "black", 16, -1,
     "第六页“Орлуулагчид ойр талаас дугаарлаж, салаа гинжид бага байрлалын дугаар өгнө.\nДугаарлалт нь дурын биш; нэрийг аль болох тодорхой, цорын ганц болгох зорилготой.”文字不是黑色、微软雅黑、16磅"),
    (6, ["Энэ хуудас гол гинжийн сонголтыг үргэлжлүүлж, “аль талаас тоолох вэ” гэдгийг шийднэ"],
     "gray", 12.5, -1,
     "第六页“Энэ хуудас гол гинжийн сонголтыг үргэлжлүүлж, ”аль талаас тоолох вэ“ гэдгийг шийднэ”文字不是灰色、微软雅黑、12.5磅"),
    # 第七页
    (7, ["“Салаанд ойр”-ыг үргэлжлүүлж, хоёр талын зай ижил онцгой тохиолдлыг шийднэ"],
     "gray", 12.5, -1,
     "第七页“”Салаанд ойр“-ыг үргэлжлүүлж, хоёр талын зай ижил онцгой тохиолдлыг шийднэ”文字不是灰色、微软雅黑、12.5磅"),
    (7, ["1. Эхлээд орлуулагчийн зай ижил эсэхийг харна 2. Зай ижил бол салааны төвөгшлийг харьцуулна 3. Энгийн салаанд бага дугаар өгнө 4. Дараа нь бүтэн нэрийг бичнэ"],
     "black", 17, -1,
     "第七页“1. Эхлээд орлуулагчийн зай ижил эсэхийг харна 2. Зай ижил бол салааны төвөгшлийг харьцуулна 3. Энгийн салаанд бага дугаар өгнө 4. Дараа нь бүтэн нэрийг бичнэ”文字不是黑色、微软雅黑、17磅"),
    (7, ["Метил нь ихэвчлэн этилаас энгийн. Хэрэв хоёр чигийн зай ижил бол метил байрласан талаас эхэлж болно.Ингэснээр нэр нь системчилсэн нэршлийн давуу дараалалтай илүү нийцнэ."],
     "black", 16, -1,
     "第七页“Метил нь ихэвчлэн этилаас энгийн. Хэрэв хоёр чигийн зай ижил бол метил байрласан талаас эхэлж болно.\nИнгэснээр нэр нь системчилсэн нэршлийн давуу дараалалтай илүү нийцнэ.”文字不是黑色、微软雅黑、16磅"),
    # 第八页
    (8, ["Салааны байрлалын хослол: 2、3、5Нийлбэр бага тул түрүүлж авна."],
     "black", 19, -1,
     "第八页“Салааны байрлалын хослол: 2、3、5Нийлбэр бага тул түрүүлж авна.”文字不是黑色、微软雅黑、19磅"),
    (8, ["Ойр → Энгийн → Бага Эхлээд ойр, дараа нь энгийн, эцэст нь байрлалын хослолыг харьцуулна."],
     "black", 20, -1,
     "第八页“Ойр → Энгийн → Бага Эхлээд ойр, дараа нь энгийн, эцэст нь байрлалын хослолыг харьцуулна.”文字不是黑色、微软雅黑、20磅"),
    (8, ["Энэ хуудас дугаарлах дүрмийг үргэлжлүүлж, “байрлалын нийлбэр бага” зарчмаар нэрийн хоёрдмол утгаас зайлсхийдэг"],
     "gray", 12.5, -1,
     "第八页“Энэ хуудас дугаарлах дүрмийг үргэлжлүүлж, ”байрлалын нийлбэр бага“ зарчмаар нэрийн хоёрдмол утгаас зайлсхийдэг”文字不是灰色、微软雅黑、12.5磅"),
    # 第九页
    (9, ["Байрлал-орлуулагчийн нэр + гол гинжийн нэрОлон өөр салаа: энгийн салаа өмнө, төвөгтэй салаа дараа.Нэрээс ганц бүтцийг буцаан гаргаж чаддаг байх ёстой."],
     "black", 18, -1,
     "第九页“Байрлал-орлуулагчийн нэр + гол гинжийн нэр\nОлон өөр салаа: энгийн салаа өмнө, төвөгтэй салаа дараа.\nНэрээс ганц бүтцийг буцаан гаргаж чаддаг байх ёстой.”文字不是黑色、微软雅黑、18磅"),
    (9, ["Энэ хуудас өмнөх гурван алхмыг бүрэн бичих хэлбэр болгон нэгтгэнэ."],
     "gray", 12.5, -1,
     "第九页“Энэ хуудас өмнөх гурван алхмыг бүрэн бичих хэлбэр болгон нэгтгэнэ.г”文字不是灰色、微软雅黑、12.5磅"),
    # 第十页
    (10, ["Энэ хуудас нэршлээс бичилт рүү шилжиж, “C хасах — шилжүүлэх — давхардлыг арилгах” аргыг хэрэглэнэ"],
     "gray", 12.5, -1,
     "第十页“Энэ хуудас нэршлээс бичилт рүү шилжиж, ”C хасах — шилжүүлэх — давхардлыг арилгах“ аргыг хэрэглэнэ”文字不是灰色、微软雅黑、12.5磅"),
    (10, ["1-р алхам: зургаан C тасралтгүй холбогдсон шулуун гинжийг бичнэ. 2-р алхам: гол гинжээс нэг C-г “салгаж”, метил болгон өөр байрлалд тавина. 3-р алхам: зүүн-баруун тэгш хэмийг анхаарч, эквивалент бүтцийг давтахгүй."],
     "black", 18, -1,
     "第十页“1-р алхам: зургаан C тасралтгүй холбогдсон шулуун гинжийг бичнэ. 2-р алхам: гол гинжээс нэг C-г ”салгаж“, метил болгон өөр байрлалд тавина. 3-р алхам: зүүн-баруун тэгш хэмийг анхаарч, эквивалент бүтцийг давтахгүй.”文字不是黑色、微软雅黑、18磅"),
    # 第十一页
    (11, ["Өмнөх хуудасны шулуун гинжин араг ясыг үргэлжлүүлж, гол гинжээс хоёр C авах тохиолдлыг авч үзнэ."],
     "gray", 12.5, -1,
     "第十一页“Өмнөх хуудасны шулуун гинжин араг ясыг үргэлжлүүлж, гол гинжээс хоёр C авах тохиолдлыг авч үзнэ.”文字不是灰色、微软雅黑、12.5磅"),
    (11, ["Хоёр C авч этил үүсгэвэл тавьсны дараа илүү урт гол гинжтэй эквивалент болохгүйг батална.",
          "Хоёр метилийг нэг C дээр, хөрш C дээр эсвэл зайтай C дээр байрлуулж үзээд, тэгш хэмээр давхардлыг арилгана.",
          "Салгасан C-ийн тоо үлдээсэн гол гинжийн C тооноос их байж болохгүй; бүтэц бүрийг бичихдээ гол гинж хамгийн урт хэвээр эсэхийг шалгана."],
     "black", 16, -1,
     "第十一页“Хоёр C авч этил үүсгэвэл тавьсны дараа илүү урт гол гинжтэй эквивалент болохгүйг батална.”、“Хоёр метилийг нэг C дээр, хөрш C дээр эсвэл зайтай C дээр байрлуулж үзээд, тэгш хэмээр давхардлыг арилгана.”、“Салгасан C-ийн тоо үлдээсэн гол гинжийн C тооноос их байж болохгүй; бүтэц бүрийг бичихдээ гол гинж хамгийн урт хэвээр эсэхийг шалгана.”文字不是黑色、微软雅黑、16磅"),
    # 第十二页
    (12, ["Энэ хуудас салаа шилжүүлэхийг үргэлжлүүлж, бүтцийн томьёоны жагсаалтыг гүйцээн шалгана"],
     "gray", 12.5, -1,
     "第十二页“Энэ хуудас салаа шилжүүлэхийг үргэлжлүүлж, бүтцийн томьёоны жагсаалтыг гүйцээн шалгана”文字不是灰色、微软雅黑、12.5磅"),
    (12, ["C атом ихэвчлэн дөрвөн ковалент холбоо үүсгэдэг. C араг ясыг зурсны дараа C бүрийн одоогийн холбооны тоонд үндэслэн H атомыг нөхөж, хураангуй бүтцийн томьёо болгон бичнэ.",
          "n-гексан 2-метилпентан 3-метилпентан 2,2-диметилбутан 2,3-диметилбутан"],
     "black", 21, -1,
     "第十二页“C атом ихэвчлэн дөрвөн ковалент холбоо үүсгэдэг. C араг ясыг зурсны дараа C бүрийн одоогийн холбооны тоонд үндэслэн H атомыг нөхөж, хураангуй бүтцийн томьёо болгон бичнэ.”、“n-гексан 2-метилпентан 3-метилпентан 2,2-диметилбутан 2,3-диметилбутан”文字不是黑色、微软雅黑、21磅"),
    # 第十三页
    (13, ["Сүүлийн хуудас өмнөх дүрмүүдтэй хаалттай холбоо үүсгэж, бодлогоор гол гинж, дугаарлалт, изомер бичилтийг бататгана"],
     "gray", 12.5, -1,
     "第十三页“Сүүлийн хуудас өмнөх дүрмүүдтэй хаалттай холбоо үүсгэж, бодлогоор гол гинж, дугаарлалт, изомер бичилтийг бататгана”文字不是灰色、微软雅黑、12.5磅"),
    (13, ["Гол гинж нь дөрвөн C-тэй, нийт C тоо зургаа болох гинжин алкан хэдэн төрлийн C араг ястай байж болох вэ?",
          "Нэг гинжин алканы харьцангуй молекул масс ойролцоогоор 100 бөгөөд молекулд гурван метил байна. Нөхцөл хангах хураангуй бүтцийн томьёог бичиж, системчилсэн нэршлээр нэрлэ."],
     "black", 15.5, -1,
     "第十三页“Гол гинж нь дөрвөн C-тэй, нийт C тоо зургаа болох гинжин алкан хэдэн төрлийн C араг ястай байж болох вэ?”、“Нэг гинжин алканы харьцангуй молекул масс ойролцоогоор 100 бөгөөд молекулд гурван метил байна. Нөхцөл хангах хураангуй бүтцийн томьёог бичиж, системчилсэн нэршлээр нэрлэ.”文字不是黑色、微软雅黑、15.5磅"),
    (13, ["Эхлээд ерөнхий томьёогоор C тоог гаргаж, дараа нь хамгийн урт гол гинжийг тогтоон, эцэст нь метилийн тоо ба дугаарлах дүрмийг шалга."],
     "black", 17, -1,
     "第十三页“Эхлээд ерөнхий томьёогоор C тоог гаргаж, дараа нь хамгийн урт гол гинжийг тогтоон, эцэст нь метилийн тоо ба дугаарлах дүрмийг шалга.”文字不是黑色、微软雅黑、17磅"),
]


# --------------------------------------------------------------------------
# 维度1：可用与可修改性（一票否决）
# --------------------------------------------------------------------------
def check_dimension1(filename: str) -> list[str]:
    fails: list[str] = []

    # 交付文件为.pptx格式，能够正常打开
    if not filename.lower().endswith(".pptx"):
        fails.append("文件不是 .pptx 格式")

    return fails


# --------------------------------------------------------------------------
# 维度2：完成度
# --------------------------------------------------------------------------
def find_shapes_with_fragment(slide, fragments: list[str]) -> list[object]:
    """返回该页中文本（去空白后）包含任一片段的形状列表。"""
    hits: list[object] = []
    frags = [norm(f) for f in fragments]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = norm(shape.text_frame.text)
        if not t:
            continue
        if any(fr and fr in t for fr in frags):
            hits.append(shape)
    return hits


def check_dimension2(prs) -> tuple[int, list[dict[str, object]]]:
    """维度2：逐条产出 Dim2Item。
    每条 rubric（得分点/扣分点）无论是否命中，都输出一条 item：
      - 命中：hit=True，delta=该条分值
      - 未命中：hit=False，delta=0
    max_delta 恒为该条 rubric 的分值上限/下限（得分点=+3，扣分点=负值）。
    """
    items: list[dict[str, object]] = []
    total: int = 0

    # ---- 得分点：+3 每一页正文文本皆为蒙古语（共 13 项，逐页产出）----
    for idx, slide in enumerate(prs.slides, 1):
        body_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                if t and t.strip():
                    body_texts.append(t)
        # "正文文本皆为蒙古语"：该页所有"正文"文本均判定为蒙古语。
        # 不计入：纯符号/数字（页码 01、02…）与化学式/化学符号（CH3、—CnH2n+1、C 等），
        # 这些不属于自然语言正文，不影响"正文皆为蒙古语"的判定。
        textual = [t for t in body_texts
                   if any(ch.isalpha() for ch in t) and not is_chemical_or_symbol(t)]
        # 使用严格判定：排除 CJK/希腊/阿拉伯/韩文等其它书写系统，
        # 剔除化学元素类拉丁字母后西里尔占比需 >= 90%，且西里尔字符数 >= 2。
        non_mong = [t for t in textual if not is_mongolian_strict(t)]
        all_mong = bool(textual) and not non_mong
        rule_desc = "+3：第%d页正文文本皆为蒙古语" % idx
        if all_mong:
            total += 3
            items.append({
                "rule": rule_desc,
                "max_delta": 3,
                "delta": 3,
                "hit": True,
                "detail": "该页正文文本皆为蒙古语",
            })
        else:
            if not textual:
                fail_detail = "该页未检出可判定的正文文本"
            else:
                sample = non_mong[0].strip().replace("\n", " ")
                if len(sample) > 40:
                    sample = sample[:40] + "…"
                fail_detail = "该页存在非蒙古语正文（共 %d 条，示例：%s）" % (
                    len(non_mong), sample)
            items.append({
                "rule": rule_desc,
                "max_delta": 3,
                "delta": 0,
                "hit": False,
                "detail": fail_detail,
            })

    # ---- 扣分点：字体/字号/颜色规则（逐条产出）----
    for (page, frags, req_color, req_size, penalty, desc) in RULES:
        rule_desc = "%d：%s" % (penalty, desc)
        if page > len(prs.slides):
            # 页码不存在 -> 视为未命中该扣分项（对应页缺失已被 -5 规则覆盖）
            items.append({
                "rule": rule_desc,
                "max_delta": penalty,
                "delta": 0,
                "hit": False,
                "detail": "PPT 不足 %d 页，无法定位该条 rubric" % page,
            })
            continue

        slide = prs.slides[page - 1]
        shapes = find_shapes_with_fragment(slide, frags)
        if not shapes:
            # 找不到目标文本 -> 视为内容已被改动/缺失，按"不符合"扣分
            total += penalty
            items.append({
                "rule": rule_desc,
                "max_delta": penalty,
                "delta": penalty,
                "hit": True,
                "detail": "未找到目标文本",
            })
            continue

        violated = False
        reasons = []          # 具体不达标项（去重，保留首次出现顺序）
        for shape in shapes:
            # 文本框开了 normAutofit -> 办公软件实际显示字号被自动放大/缩小，
            # 与细则存储字号不符（本文件即用此手法把 13.5 显示成 18 等），直接判字号不达标。
            autofit = shape_has_autofit(shape)
            for ri in shape_runs(shape):
                bad = []
                # 字体
                if ri.font != FONT_REQ:
                    bad.append("字体不是%s（实际：%s）" % (FONT_REQ, ri.font or "未知"))
                # 颜色色系
                if ri.color_cls != req_color:
                    cn = {"black": "黑色", "gray": "灰色",
                          "darkgreen": "深绿色", "green": "绿色"}
                    bad.append("颜色不是%s（实际：%s）" % (
                        cn.get(req_color, req_color),
                        cn.get(ri.color_cls, ri.color_cls) if ri.color_cls else "未知"))
                # 字号：autofit 开启时显示字号不可信，直接判不达标；否则严格比对存储字号
                if autofit:
                    bad.append("字号不达标（文本框开启自动调整，办公软件实际显示≠%s磅）" % req_size)
                elif not size_matches(ri.size_pt, req_size):
                    bad.append("字号不是%s磅（实际：%s磅）" % (
                        req_size,
                        ("%g" % ri.size_pt) if ri.size_pt is not None else "未知"))
                if bad:
                    violated = True
                    for b in bad:
                        if b not in reasons:
                            reasons.append(b)

        if violated:
            total += penalty
            loc = desc.split("”")[0] + "”" if "”" in desc else desc
            items.append({
                "rule": rule_desc,
                "max_delta": penalty,
                "delta": penalty,
                "hit": True,
                "detail": "%s —— %s" % (loc, "；".join(reasons)),
            })
        else:
            items.append({
                "rule": rule_desc,
                "max_delta": penalty,
                "delta": 0,
                "hit": False,
                "detail": "目标文本存在且字体/颜色/字号均符合要求",
            })

    # ---- 扣分点：-5 PPT页数不是13页 ----
    page_count = len(prs.slides)
    if page_count != 13:
        total += -5
        items.append({
            "rule": "-5：PPT页数不是13页",
            "max_delta": -5,
            "delta": -5,
            "hit": True,
            "detail": "当前 %d 页" % page_count,
        })
    else:
        items.append({
            "rule": "-5：PPT页数不是13页",
            "max_delta": -5,
            "delta": 0,
            "hit": False,
            "detail": "PPT 共 13 页",
        })

    # ---- 扣分点：-3 PPT出现了文字遮挡情况或文字显示不全 ----
    #   细则两种情形：
    #     (a) 文字遮挡：一段文字被别的东西盖住一部分（一半可见、一半被遮）——
    #         表现为两个"都承载实质正文文字"的形状在页面上部分重叠，
    #         互相压住对方一部分（互不完全包含）。
    #         注：把小标签/化学式(CH3、—R 等)有意放在大文本框留白处属正常图文叠放，不算遮挡。
    #     (b) 文字显示不全：文本框内容放不下被裁切——
    #         normAutofit 自动缩字(fontScale<100%) 或 文本形状超出幻灯片边界。
    sw, sh_ = prs.slide_width, prs.slide_height

    def body_text_shapes(slide):
        """承载"实质正文文字"（含蒙文、且非纯化学式/符号）的文本形状。"""
        out = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text
            if not txt.strip():
                continue
            if is_chemical_or_symbol(txt) or not is_mongolian(txt):
                continue
            try:
                l, t, w, h = shape.left, shape.top, shape.width, shape.height
            except Exception:
                continue
            if None in (l, t, w, h):
                continue
            out.append((shape, l, t, w, h))
        return out

    def all_text_shapes(slide):
        out = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                try:
                    l, t, w, h = shape.left, shape.top, shape.width, shape.height
                except Exception:
                    continue
                if None in (l, t, w, h):
                    continue
                out.append((shape, l, t, w, h))
        return out

    def overlap_area(a, b):
        _, l1, t1, w1, h1 = a
        _, l2, t2, w2, h2 = b
        ix = max(0, min(l1 + w1, l2 + w2) - max(l1, l2))
        iy = max(0, min(t1 + h1, t2 + h2) - max(t1, t2))
        return ix * iy

    # 相邻/相接文本框常因排版仅共享 1~2 EMU 边缘，属正常贴边，不算遮挡。
    # 仅当重叠面积超过较小文本框面积的一定比例（10%）才视为"真正盖住了一部分文字"。
    OCC_RATIO = 0.10

    occluded_pages = []     # (a) 文字遮挡
    clipped_pages = []      # (b) 文字显示不全
    for idx, slide in enumerate(prs.slides, 1):
        bts = body_text_shapes(slide)
        occ = False
        for i in range(len(bts)):
            for j in range(i + 1, len(bts)):
                ov = overlap_area(bts[i], bts[j])
                if ov <= 0:
                    continue
                small = max(1, min(bts[i][3] * bts[i][4], bts[j][3] * bts[j][4]))
                if ov / small < OCC_RATIO:
                    continue
                occ = True
                break
            if occ:
                break
        if occ:
            occluded_pages.append(idx)
        clip = False
        for shape, l, t, w, h in all_text_shapes(slide):
            if shape_text_overflow(shape):
                clip = True
                break
            if l < 0 or t < 0 or (l + w) > sw or (t + h) > sh_:
                clip = True
                break
        if clip:
            clipped_pages.append(idx)

    if occluded_pages or clipped_pages:
        total += -3
        detail = []
        if occluded_pages:
            detail.append("文字遮挡：第%s页" % "、".join(map(str, occluded_pages)))
        if clipped_pages:
            detail.append("文字显示不全：第%s页" % "、".join(map(str, clipped_pages)))
        items.append({
            "rule": "-3：PPT出现了文字遮挡情况或文字显示不全",
            "max_delta": -3,
            "delta": -3,
            "hit": True,
            "detail": "；".join(detail),
        })
    else:
        items.append({
            "rule": "-3：PPT出现了文字遮挡情况或文字显示不全",
            "max_delta": -3,
            "delta": 0,
            "hit": False,
            "detail": "未检出文字遮挡或显示不全",
        })

    return total, items


# --------------------------------------------------------------------------
# 统一入口
# --------------------------------------------------------------------------
def _locate_pptx(dir_path: str) -> str | None:
    """在给定目录中定位待评估的 .pptx 文件。
    优先选取非隐藏、非临时（不以 ~$ 开头）的 Office 文件。"""
    if not os.path.isdir(dir_path):
        return None
    candidates: list[str] = []
    for name in os.listdir(dir_path):
        if name.startswith("~$"):
            continue
        low = name.lower()
        if low.endswith(".pptx"):
            candidates.append(name)
    if not candidates:
        return None
    # 按文件名排序，保持稳定
    candidates.sort()
    return os.path.join(dir_path, candidates[0])


def evaluate(dir_path: str) -> dict:
    """统一入口：接收脚本所在目录的路径，在该目录里定位并评估 PPT 文档。"""
    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": 3 * 13,  # 得分点满分（每页 +3，共 13 页）
    }

    try:
        file_path = _locate_pptx(dir_path)
        if not file_path:
            result["status"] = "error"
            result["error"] = "目录中未找到 .pptx 文件"
            return result

        filename: str = os.path.basename(file_path)
        result["file_name"] = filename

        try:
            prs = Presentation(file_path)
        except Exception as e:
            result["dim1_pass"] = False
            result["dim1_reason"] = "文件无法打开/解析：%s" % e
            return result

        # 维度1
        d1_fails = check_dimension1(filename)
        if d1_fails:
            result["dim1_pass"] = False
            result["dim1_reason"] = "；".join(d1_fails)
            return result

        result["dim1_pass"] = True

        # 维度2：得分点 + 扣分点（每条 rubric 均产出一条 item，未命中时 hit=False、delta=0）
        total, dim2_items = check_dimension2(prs)
        result["dim2_items"] = dim2_items
        result["total_score"] = total
    except Exception as e:
        result["status"] = "error"
        result["error"] = "%s" % e

    return result


if __name__ == "__main__":
    import json
    target_dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(
        os.path.abspath(__file__))
    print(json.dumps(evaluate(target_dir), ensure_ascii=False, indent=2))
