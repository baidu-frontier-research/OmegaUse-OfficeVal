# -*- coding: utf-8 -*-
"""
PPT Auto Evaluation Script
Target: 幼儿游戏活动课程设计与实施_修改完成_含音乐.pptx

统一约定:
- 仅暴露 evaluate(file_path: str) -> dict
- 不修改 sys.stdout, 不 sys.exit, 主结果只通过 return 输出
- 文档路径全部由参数传入, 不硬编码路径
"""
import sys, os, re, json, zipfile

SCRIPT_ID = "042"
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE


# 由 evaluate(file_path) 在调用时设置,供内部 helper 读取包内资源使用。
# 不做任何默认路径假设,避免模块级路径写死。
_FILE_PATH = None


def get_shape_xml(shape):
    return shape._element.xml

def shape_has_image_fill(shape):
    xml = get_shape_xml(shape)
    return 'blipFill' in xml or 'r:embed' in xml

def get_blip_alpha(shape):
    # 返回图片透明度百分比。OOXML中alphaModFix amt为"不透明度"(千分比),
    # 办公软件里的"透明度" = 100 - 不透明度
    xml = get_shape_xml(shape)
    m = re.search(r'alphaModFix.*?amt="(\d+)"', xml)
    if m:
        return 100 - int(m.group(1)) / 1000.0
    m = re.search(r'<a:blip[^/]*>(.*?)</a:blip>', xml, re.DOTALL)
    if m:
        am = re.search(r'alphaModFix.*?amt="(\d+)"', m.group(1))
        if am:
            return 100 - int(am.group(1)) / 1000.0
    return None

def get_shape_fill_alpha(shape):
    # 返回形状图片填充的透明度百分比(同上,透明度 = 100 - 不透明度)
    xml = get_shape_xml(shape)
    blip_match = re.search(r'<a:blip[^>]*>(.*?)</a:blip>', xml, re.DOTALL)
    if blip_match:
        am = re.search(r'alphaModFix.*?amt="(\d+)"', blip_match.group(1))
        if am:
            return 100 - int(am.group(1)) / 1000.0
    return None

def run_uses_arial(run):
    xml = run._r.xml
    typefaces = dict(re.findall(r'<a:(latin|ea|cs)[^>]*typeface="([^"]*)"', xml))
    # 与办公软件字体框保持一致:软件按run的语言(lang)报告字体框显示的字体。
    # lang为英文(如en-US)时,字体框显示的是西文字体(latin),肉眼可见即为该字体;
    # 此时应按latin判定,而非强制读中文字体(ea)。
    lang_m = re.search(r'\blang="([^"]*)"', xml)
    lang = (lang_m.group(1) if lang_m else '').lower()
    if lang.startswith('en'):
        font_name = typefaces.get('latin') or run.font.name or ''
    elif any('一' <= ch <= '鿿' for ch in run.text):
        font_name = typefaces.get('ea') or run.font.name or ''
    else:
        font_name = typefaces.get('latin') or run.font.name or ''
    return 'arial' in font_name.lower()

_theme_clr_cache = {}

def _load_theme_colors(file_path):
    # 从主题(theme1.xml)读取配色方案,把主题色名映射为实际RGB。
    # 这样schemeClr(如accent1)也能解析成肉眼可见的真实颜色。
    if file_path in _theme_clr_cache:
        return _theme_clr_cache[file_path]
    colors = {}
    try:
        with zipfile.ZipFile(file_path) as z:
            xml = z.read('ppt/theme/theme1.xml').decode('utf-8')
        block = re.search(r'<a:clrScheme.*?</a:clrScheme>', xml, re.DOTALL)
        if block:
            for name in ('dk1','lt1','dk2','lt2','accent1','accent2',
                         'accent3','accent4','accent5','accent6','hlink','folHlink'):
                mm = re.search(r'<a:'+name+r'>(.*?)</a:'+name+r'>', block.group(0), re.DOTALL)
                if mm:
                    c = (re.search(r'val="([0-9A-Fa-f]{6})"', mm.group(1))
                         or re.search(r'lastClr="([0-9A-Fa-f]{6})"', mm.group(1)))
                    if c:
                        colors[name] = c.group(1).upper()
        # schemeClr里常用别名(dk1/lt1对应窗口文本/背景色)
        colors.setdefault('tx1', colors.get('dk1', ''))
        colors.setdefault('bg1', colors.get('lt1', ''))
        colors.setdefault('tx2', colors.get('dk2', ''))
        colors.setdefault('bg2', colors.get('lt2', ''))
    except Exception:
        pass
    _theme_clr_cache[file_path] = colors
    return colors

def get_run_rgb(run, file_path):
    # 获取run字体颜色的RGB十六进制字符串。只要颜色符合即可,不论来源:
    # 直接RGB(srgbClr)直接返回;主题色(schemeClr)则解析主题得到真实RGB。
    try:
        color = run.font.color
        if color and color.type is not None:
            # 直接RGB色
            try:
                return str(color.rgb).upper()
            except Exception:
                pass
            # 主题色: 解析run XML里的schemeClr名,映射到主题实际RGB
            xml = run._r.xml
            sm = re.search(r'<a:solidFill>\s*<a:schemeClr val="([^"]+)"', xml)
            if sm:
                return _load_theme_colors(file_path).get(sm.group(1), '')
    except Exception:
        pass
    return ''


def emu_to_cm(emu):
    return emu / 360000.0

def get_all_text(slide):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if p.text.strip():
                    texts.append(p.text.strip())
    return texts


def _shape_text(shape):
    if not getattr(shape, 'has_text_frame', False):
        return ''
    return '\n'.join(p.text for p in shape.text_frame.paragraphs).strip()


def _find_labeled_shape_index(slide, label):
    """按独立文本行定位标签，兼容编号与标题合并在同一文本框。"""
    for index, shape in enumerate(slide.shapes):
        text = _shape_text(shape)
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        if label in lines or text == label:
            return index
    return None


# 维度二所有评分项(命中/未命中均需返回),max_delta 为该项满分,
# 命中时 delta = max_delta(加分项)或负值(扣分项),未命中时 delta = 0。
_DIM2_RULES = [
    ("第1页背景替换为图片,高18-20cm宽32-34cm,无边缘留白或空白块", 1),
    ("第1页背景图片透明度约50%(45-55%),背景不明显压过标题和正文文字", 3),
    ("第2页背景替换为图片,高18-20cm宽32-34cm,无平铺/严重裁切/边缘留白", 1),
    ("第2页背景图片透明度约85%(80-90%),背景较淡不影响页面文字识别", 3),
    ("第2页红色标题[整体教学设计]红色字体为Arial、加粗", 1),
    ("第2页六个小标题(学习诊断/任务转化/资源整合/情境演练/表现评价/迭代成长)均为蓝色字体+下划线", 1),
    ("第2页[学习诊断/任务转化/资源整合/情境演练/表现评价/迭代成长]小标题均为竖排显示", 5),
    ("第2页[学生/任务/资源]三个标题均横排,居中于对应内容框上方", 1),
    ("第2页粉/黄/蓝三栏内容框宽6-8cm、高3-4.3cm,呈三栏结构", 3),
    ("第2页未出现独立的[评价]小标题及对应内容", 1),
    ("第6页音频播放按钮在左下角,可见且不遮挡标题/正文/页码", 3),
    ("第6页音频可点击播放,且以嵌入或可靠打包方式存在", 3),
]

_DEDUCTION_RULES = [
    ("第2页整体为一张图片(无可编辑文字)", -5),
    ("第3页教学目标及下方文本被删除,或被移动到第4页标题之后", -3),
    ("第3页标题位于页面外", -1),
    ("第3页品格目标与知识目标重叠", -1),
    ("除封面外任意页页码丢失", -1),
    ("第1页或第2页背景为0%的透明度", -1),
    ("第1页标题[幼儿游戏活动 课程设计与实施]未位于页面左侧", -1),
]


# 目标文档所在目录内待评估文件的期望名称(脚本自扫描,不依赖调用方传入具体文件)
_TARGET_NAME = "幼儿游戏活动课程设计与实施_修改完成_含音乐.pptx"


def _resolve_file(dir_path: str):
    # 在 dir_path 目录内定位被评估文档: 优先匹配期望文件名,
    # 若不存在再回退为该目录下第一个 .pptx 文件(忽略临时锁文件"~$*")。
    # 返回 (file_path, error_msg);找不到时 file_path 为 None。
    if not dir_path or not os.path.isdir(dir_path):
        return None, f"目录不存在或不是目录: {dir_path}"
    expected = os.path.join(dir_path, _TARGET_NAME)
    if os.path.isfile(expected):
        return expected, None
    candidates = []
    for name in os.listdir(dir_path):
        if name.startswith('~$'):
            continue
        ext = os.path.splitext(name)[1].lower()
        if ext == '.pptx' and os.path.isfile(os.path.join(dir_path, name)):
            candidates.append(name)
    if not candidates:
        return None, f"目录内未找到 .pptx 文件: {dir_path}"
    # 多个候选时优先取名称与期望文件最接近的(简单以最长公共前缀为序)
    candidates.sort(key=lambda n: -len(os.path.commonprefix([n, _TARGET_NAME])))
    return os.path.join(dir_path, candidates[0]), None


def evaluate(dir_path: str) -> dict:
    global _FILE_PATH
    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": sum(md for _, md in _DIM2_RULES if md > 0),
    }

    # 维度二初始化: 全部记为未命中(命中/未命中均返回,便于批量汇总对齐)
    dim2 = [
        {"rule": r, "max_delta": md, "delta": 0, "hit": False, "detail": ""}
        for r, md in _DIM2_RULES
    ]
    deductions = [
        {"rule": r, "max_delta": md, "delta": 0, "hit": False, "detail": ""}
        for r, md in _DEDUCTION_RULES
    ]

    try:
        # 在给定目录内自行定位被评估文档
        file_path, err = _resolve_file(dir_path)
        if err:
            result["status"] = "error"
            result["error"] = err
            return result
        _FILE_PATH = file_path
        result["file_name"] = os.path.basename(file_path)

        ext = os.path.splitext(file_path)[1].lower()
        if ext != '.pptx':
            result["dim1_pass"] = False
            result["dim1_reason"] = "文件格式不是.pptx"
            return result
        try:
            prs = Presentation(file_path)
        except Exception as e:
            result["dim1_pass"] = False
            result["dim1_reason"] = f"文件无法正常打开: {e}"
            return result

        if len(prs.slides) < 6:
            result["dim1_reason"] = f"PPT 页数不足 6 页，实际 {len(prs.slides)} 页"
            return result

        slide1, slide2 = prs.slides[0], prs.slides[1]

        # 维度一全部通过
        result["dim1_pass"] = True

        # === 维度二 ===

        # +1: 第1页背景替换为图片,高18-20cm 宽32-34cm,没有边缘留白或空白块
        slide1_bg = None
        for s in slide1.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if emu_to_cm(s.width) > 25 and emu_to_cm(s.height) > 15:
                    slide1_bg = s; break
        if slide1_bg:
            w, h = emu_to_cm(slide1_bg.width), emu_to_cm(slide1_bg.height)
            size_ok = 32 <= w <= 34 and 18 <= h <= 20
            # 没有边缘留白或空白块: 图片须覆盖整页四条边缘(办公软件中不露白边)
            tol = Cm(0.5)
            no_edge_blank = (
                slide1_bg.left <= tol and slide1_bg.top <= tol
                and slide1_bg.left + slide1_bg.width >= prs.slide_width - tol
                and slide1_bg.top + slide1_bg.height >= prs.slide_height - tol
            )
            if size_ok and no_edge_blank:
                dim2[0]["hit"] = True
                dim2[0]["delta"] = 1
                dim2[0]["detail"] = ""
            elif not size_ok:
                dim2[0]["detail"] = ""
            else:
                dim2[0]["detail"] = ""
        else:
            dim2[0]["detail"] = ""

        # +3: 第1页背景图片透明度约50%(45-55%),背景不明显压过标题和正文文字
        if slide1_bg:
            alpha = get_blip_alpha(slide1_bg)
            alpha_ok = alpha is not None and 45 <= alpha <= 55
            # 不压过文字: 背景图片须位于所有正文/标题文字之下(z序在底层),
            # 办公软件按形状顺序渲染,靠前的在下层,文字才不会被背景覆盖
            bg_idx = list(slide1.shapes).index(slide1_bg)
            text_idxs = [i for i, s in enumerate(slide1.shapes)
                         if s.has_text_frame and any(p.text.strip() for p in s.text_frame.paragraphs)]
            not_cover_text = all(bg_idx < ti for ti in text_idxs) if text_idxs else True
            if alpha_ok and not_cover_text:
                dim2[1]["hit"] = True
                dim2[1]["delta"] = 3
                dim2[1]["detail"] = ""
            elif not alpha_ok:
                dim2[1]["detail"] = ""
            else:
                dim2[1]["detail"] = ""
        else:
            dim2[1]["detail"] = ""

        # +1: 第2页背景替换为图片,高18-20cm宽32-34cm,无平铺/严重裁切/边缘留白或空白块
        slide2_bg = None
        for s in slide2.shapes:
            if shape_has_image_fill(s) and emu_to_cm(s.width) > 25 and emu_to_cm(s.height) > 15:
                slide2_bg = s; break
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE and emu_to_cm(s.width) > 25 and emu_to_cm(s.height) > 15:
                slide2_bg = s; break
        if slide2_bg:
            w, h = emu_to_cm(slide2_bg.width), emu_to_cm(slide2_bg.height)
            xml = get_shape_xml(slide2_bg)
            size_ok = 32 <= w <= 34 and 18 <= h <= 20
            # 无平铺: 填充为拉伸(stretch),不含tile平铺
            has_tile = 'tile' in xml.lower() and 'stretch' not in xml.lower()
            # 无严重裁切: srcRect各边裁切比例均不超过25%(办公软件裁切以千分比记)
            severe_crop = False
            sr = re.search(r'<a:srcRect([^/>]*)/?>', xml)
            if sr:
                for side in ('l', 't', 'r', 'b'):
                    cm = re.search(side + r'="(-?\d+)"', sr.group(1))
                    if cm and abs(int(cm.group(1))) > 25000:
                        severe_crop = True
            # 无边缘留白或空白块: 图片四边覆盖整页(办公软件中不露白边)
            tol = Cm(0.5)
            no_edge_blank = (
                slide2_bg.left <= tol and slide2_bg.top <= tol
                and slide2_bg.left + slide2_bg.width >= prs.slide_width - tol
                and slide2_bg.top + slide2_bg.height >= prs.slide_height - tol
            )
            if size_ok and not has_tile and not severe_crop and no_edge_blank:
                dim2[2]["hit"] = True
                dim2[2]["delta"] = 1
                dim2[2]["detail"] = ""
            elif not size_ok:
                dim2[2]["detail"] = ""
            elif has_tile:
                dim2[2]["detail"] = ""
            elif severe_crop:
                dim2[2]["detail"] = ""
            else:
                dim2[2]["detail"] = ""
        else:
            dim2[2]["detail"] = ""

        # +3: 第2页背景图片透明度约85%(80-90%),背景较淡不影响页面文字识别
        if slide2_bg:
            alpha = get_shape_fill_alpha(slide2_bg)
            alpha_ok = alpha is not None and 80 <= alpha <= 90
            # 不影响文字识别: 背景须位于所有正文/标题文字之下(z序在底层),
            # 办公软件按形状顺序渲染,背景在文字下层文字才清晰可辨
            bg_idx = list(slide2.shapes).index(slide2_bg)
            text_idxs = [i for i, s in enumerate(slide2.shapes)
                         if s.has_text_frame and any(p.text.strip() for p in s.text_frame.paragraphs)]
            not_affect_text = all(bg_idx < ti for ti in text_idxs) if text_idxs else True
            if alpha_ok and not_affect_text:
                dim2[3]["hit"] = True
                dim2[3]["delta"] = 3
                dim2[3]["detail"] = ""
            elif not alpha_ok:
                dim2[3]["detail"] = ""
            else:
                dim2[3]["detail"] = ""
        else:
            dim2[3]["detail"] = ""

        # +1: 第2页红色标题"整体教学设计:从学习基础到实践改进"等红色字体为Arial、加粗
        # 按文本定位标题(不再固定索引);颜色采用红色系范围判定,涵盖纯红FF0000、
        # 深红C00000、主题红(schemeClr)等办公软件常见红色;并覆盖"等红色字体":
        # 第2页内所有红色字体均须为Arial加粗。
        TITLE_KW = '整体教学设计'

        def is_red_rgb(rgb_str):
            # 红色判定: R分量为主导且显著高于G/B, 允许一定的近似范围
            # (办公软件中红色系颜色的通用特征)
            try:
                r = int(rgb_str[0:2], 16); g = int(rgb_str[2:4], 16); b = int(rgb_str[4:6], 16)
            except (ValueError, IndexError):
                return False
            return r >= 128 and r > g + 40 and r > b + 40

        title_found = False
        title_has_red = False
        all_red_ok = True
        for s in slide2.shapes:
            if not s.has_text_frame:
                continue
            is_title_shape = any(TITLE_KW in p.text for p in s.text_frame.paragraphs)
            if is_title_shape:
                title_found = True
            for p in s.text_frame.paragraphs:
                for run in p.runs:
                    if not run.text.strip():
                        continue
                    # 仅针对红色字体检查: 判定该run是否为红色系
                    if not is_red_rgb(get_run_rgb(run, file_path)):
                        continue
                    if is_title_shape:
                        title_has_red = True
                    # 红色字体须为Arial
                    if not run_uses_arial(run):
                        all_red_ok = False
                    # 红色字体须加粗
                    if not run.font.bold:
                        all_red_ok = False
        # 需同时满足: 标题按文本定位到, 标题内确有红色字体, 且页内所有红色字体均为Arial加粗
        if title_found and title_has_red and all_red_ok:
            dim2[4]["hit"] = True
            dim2[4]["delta"] = 1
            dim2[4]["detail"] = ""
        else:
            dim2[4]["detail"] = ""

        # +1: 第2页内六个小标题"学习诊断、任务转化、资源整合、情境演练、表现评价、迭代成长"
        #     均改为蓝色字体并添加下划线
        sub_titles = ['学习诊断', '任务转化', '资源整合', '情境演练', '表现评价', '迭代成长']
        sub_indices = [
            _find_labeled_shape_index(slide2, name) for name in sub_titles
        ]
        sub_shapes = [
            slide2.shapes[index]
            for index in sub_indices
            if index is not None
        ]

        def is_blue_rgb(rgb_str):
            # 蓝色判定: 蓝色分量占主导(办公软件中蓝色系颜色的通用特征)
            try:
                r = int(rgb_str[0:2], 16); g = int(rgb_str[2:4], 16); b = int(rgb_str[4:6], 16)
            except (ValueError, IndexError):
                return False
            return b >= 128 and b > r and b >= g

        all_blue_ul = len(sub_shapes) == len(sub_titles)
        for shape, name in zip(sub_shapes, sub_titles):
            # 该小标题确须存在(文本匹配细则指定的六个标题)
            if name not in _shape_text(shape):
                all_blue_ul = False
                break
            # 该小标题的每个文字run均须为蓝色字体且加下划线
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    if not run.text.strip():
                        continue
                    c = get_run_rgb(run, file_path)
                    if not is_blue_rgb(c) or not run.font.underline:
                        all_blue_ul = False
        if all_blue_ul:
            dim2[5]["hit"] = True
            dim2[5]["delta"] = 1
            dim2[5]["detail"] = ""
        else:
            dim2[5]["detail"] = ""

        # +5: 第2页所有小标题均为竖排显示
        # 办公软件竖排文字对应bodyPr的vert属性为竖排取值(vert/vert270/eaVert/mongolianVert等),
        # 缺省或horz为横排
        vertical_vals = ('vert', 'vert270', 'eaVert', 'mongolianVert', 'wordArtVert', 'wordArtVertRtl')
        all_vert = len(sub_shapes) == len(sub_titles)
        for shape in sub_shapes:
            xml = get_shape_xml(shape)
            body_m = re.search(r'<a:bodyPr([^>]*)>', xml)
            attrs = body_m.group(1) if body_m else ''
            vert_m = re.search(r'vert="([^"]+)"', attrs)
            if not vert_m or vert_m.group(1) not in vertical_vals:
                all_vert = False
                break
        if all_vert:
            dim2[6]["hit"] = True
            dim2[6]["delta"] = 5
            dim2[6]["detail"] = ""
        else:
            dim2[6]["detail"] = ""

        # +1: 第2页"学生""任务""资源"三个标题均改为横排,并分别位于对应内容框上方居中显示
        col_titles = ['学生', '任务', '资源']
        column_pairs = []
        for name in col_titles:
            title_index = _find_labeled_shape_index(slide2, name)
            if title_index is None or title_index + 1 >= len(slide2.shapes):
                continue
            column_pairs.append(
                (slide2.shapes[title_index], slide2.shapes[title_index + 1], name)
            )
        col_ok = len(column_pairs) == len(col_titles)
        for ts, cs, name in column_pairs:
            # 该标题确为"学生/任务/资源"之一
            title_text = _shape_text(ts)
            if name not in title_text:
                col_ok = False
                break
            # 均改为横排: bodyPr无vert或vert为horz(办公软件横排文字)
            xml = get_shape_xml(ts)
            bm = re.search(r'<a:bodyPr([^>]*)>', xml)
            attrs = bm.group(1) if bm else ''
            vm = re.search(r'vert="([^"]+)"', attrs)
            is_h = not vm or vm.group(1) == 'horz'
            # 位于对应内容框上方: 标题底边不低于内容框顶边(允许少量容差)
            is_above = ts.top + ts.height <= cs.top + Cm(0.5)
            # 居中显示: 标题水平中心与内容框水平中心对齐
            tc = ts.left + ts.width // 2
            cc = cs.left + cs.width // 2
            is_ctr = abs(tc - cc) < Cm(2)
            if not (is_h and is_above and is_ctr):
                col_ok = False; break
        if col_ok:
            dim2[7]["hit"] = True
            dim2[7]["delta"] = 1
            dim2[7]["detail"] = ""
        else:
            dim2[7]["detail"] = ""

        # +3: 第2页粉色、黄色、蓝色内容框重新排成三栏结构,三栏宽度6-8cm和高度3-4.3cm
        boxes = [content for _title, content, _name in column_pairs]
        ws = [emu_to_cm(b.width) for b in boxes]
        hs = [emu_to_cm(b.height) for b in boxes]
        # 三栏宽度6-8cm
        w_ok = len(boxes) == 3 and all(6 <= w <= 8 for w in ws)
        # 三栏高度3-4.3cm
        h_ok = len(boxes) == 3 and all(3 <= h <= 4.3 for h in hs)
        # 排成三栏结构: 三个框左右并排(水平位置各不相同,且大致处于同一水平行)
        lefts = sorted(b.left for b in boxes)
        is_three_col = len(lefts) == 3 and lefts[0] < lefts[1] < lefts[2]
        tops = [b.top for b in boxes]
        same_row = len(tops) == 3 and max(tops) - min(tops) < Cm(2)
        if w_ok and h_ok and is_three_col and same_row:
            dim2[8]["hit"] = True
            dim2[8]["delta"] = 3
            dim2[8]["detail"] = ""
        else:
            dim2[8]["detail"] = ""

        # +1: 第二页没有出现"评价"小标题及对应内容
        # 只针对独立的"评价"小标题(整段文字恰为"评价"),
        # 六步流程中的"表现评价"为合法内容,不在此约束
        has_eval = False
        for s in slide2.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    t = p.text.strip()
                    # 独立"评价"小标题,或以"评价"单独成栏的对应内容标题
                    if t == '评价':
                        has_eval = True
        # "评价"不应作为三栏标题出现(对应内容框也随之不存在)
        for ts, _cs, _name in column_pairs:
            if _shape_text(ts) == '评价':
                has_eval = True
        if not has_eval:
            dim2[9]["hit"] = True
            dim2[9]["delta"] = 1
            dim2[9]["detail"] = ""
        else:
            dim2[9]["detail"] = ""

        # +3: 最后一页音频播放按钮在页面左下角区域,且为PPT音频对象(非静态图片或无效占位符),
        #     按钮可见但不遮挡主要标题、正文卡片或页码
        slide6 = prs.slides[5]
        s6xml = slide6._element.xml
        # 为PPT音频对象: 含audioFile关系(真实音频对象,区别于静态图片/无效占位符)
        audio_found = 'audioFile' in s6xml
        audio_left_bottom = False
        audio_visible = False
        audio_no_overlap = False
        if audio_found:
            ps = s6xml.rfind('<p:pic>', 0, s6xml.find('audioFile'))
            pe = s6xml.find('</p:pic>', s6xml.find('audioFile'))
            if ps >= 0 and pe >= 0:
                pxml = s6xml[ps:pe+8]
                om = re.search(r'<a:off x="(-?\d+)" y="(-?\d+)"', pxml)
                em = re.search(r'<a:ext cx="(\d+)" cy="(\d+)"', pxml)
                if om and em:
                    ax, ay = int(om.group(1)), int(om.group(2))
                    aw, ah = int(em.group(1)), int(em.group(2))
                    sh = emu_to_cm(prs.slide_height)
                    sw = emu_to_cm(prs.slide_width)
                    # 位于页面左下角区域
                    audio_left_bottom = emu_to_cm(ax) < sw/3 and emu_to_cm(ay) > sh*2/3
                    # 按钮可见: 有实际尺寸且未被隐藏
                    audio_visible = aw > 0 and ah > 0 and 'hidden="1"' not in pxml
                    # 不遮挡主要标题、正文卡片或页码: 音频矩形与承载内容的形状(有文字的标题/页码、
                    # 图片/正文卡片)无实质性遮挡;背景整页图形、装饰线、音频自身图标不算内容,不参与判定。
                    # 办公软件中边角轻微相交不构成遮挡,以重叠面积占比判定实质遮挡
                    ar1, at1, ar2, at2 = ax, ay, ax + aw, ay + ah
                    audio_area = aw * ah
                    audio_no_overlap = True
                    for s in slide6.shapes:
                        if 'audioFile' in s._element.xml:
                            continue  # 跳过音频对象本身
                        if s.left is None or s.top is None or not s.width or not s.height:
                            continue
                        # 仅对内容形状判定: 有可见文字(标题/页码) 或 图片(正文卡片)
                        has_text = s.has_text_frame and any(p.text.strip() for p in s.text_frame.paragraphs)
                        is_pic = s.shape_type == MSO_SHAPE_TYPE.PICTURE
                        if not (has_text or is_pic):
                            continue
                        # 排除整页背景(近似覆盖全页的形状不属于被遮挡的标题/卡片/页码)
                        if emu_to_cm(s.width) > 25 and emu_to_cm(s.height) > 15:
                            continue
                        sl_, st_ = s.left, s.top
                        sr_, sb_ = s.left + s.width, s.top + s.height
                        ox = max(0, min(ar2, sr_) - max(ar1, sl_))
                        oy = max(0, min(at2, sb_) - max(at1, st_))
                        ov = ox * oy
                        # 实质遮挡: 重叠面积超过音频按钮或该内容形状面积的1/4
                        shape_area = s.width * s.height
                        if ov > 0 and (ov > audio_area * 0.25 or ov > shape_area * 0.25):
                            audio_no_overlap = False; break
        if audio_found and audio_left_bottom and audio_visible and audio_no_overlap:
            dim2[10]["hit"] = True
            dim2[10]["delta"] = 3
            dim2[10]["detail"] = ""
        elif audio_found and audio_left_bottom and not audio_no_overlap:
            dim2[10]["detail"] = ""
        elif audio_found and audio_left_bottom and not audio_visible:
            dim2[10]["detail"] = ""
        elif audio_found and audio_left_bottom:
            dim2[10]["detail"] = ""
        elif audio_found:
            dim2[10]["detail"] = ""
        else:
            dim2[10]["detail"] = ""

        # +3: 放映模式下点击最后一页左下角播放按钮后音乐可正常播放,
        #     且音频以嵌入或可靠打包方式存在,交付后不依赖本地绝对路径
        # 关键: 点击动作须绑定到"左下角那个音频按钮"本身,而不是仅凭页面里有
        # audioFile 就算数。做法:
        #   (a) 解析该按钮 <p:pic> 的 cNvPr 上是否带 a:hlinkClick action="ppaction://media"
        #   (b) 或在 <p:timing> 中查找 nodeType="clickEffect" 且其内引用该按钮 spid
        #       并携带媒体节点(<p:audio>)/媒体动作(ppaction://media),
        #       等价于办公软件里"点击时播放"的时间线绑定
        #   (c) 顺该按钮 <a:audioFile r:link/r:embed="rIdX"> 反查 slide6.xml.rels,
        #       确认目标为包内 ../media/ 且非 TargetMode="External"(即真正嵌入,
        #       不是指向本地绝对路径的外链)
        click_to_play = False
        audio_packaged = False
        if audio_found and ps >= 0 and pe >= 0:
            # (1) 该按钮引用的音频关系Id
            audio_rid_m = re.search(
                r'<a:audioFile\b[^/>]*r:(?:link|embed)="([^"]+)"', pxml
            )
            audio_rid = audio_rid_m.group(1) if audio_rid_m else None

            # (2a) 按钮自身的点击超链接触发媒体动作
            hlink_click_media = bool(re.search(
                r'<a:hlinkClick\b[^/>]*action="ppaction://media[^"]*"', pxml
            ))

            # (2b) 兼容: 无 hlinkClick 时,看 timing 里是否有 clickEffect 绑定到该按钮
            click_effect_bound = False
            if not hlink_click_media:
                spid_m = re.search(r'<p:cNvPr\b[^>]*\bid="([^"]+)"', pxml)
                spid = spid_m.group(1) if spid_m else None
                tm_all = re.search(r'<p:timing\b.*?</p:timing>', s6xml, re.DOTALL)
                if spid and tm_all:
                    timing_xml = tm_all.group(0)
                    for eff in re.finditer(
                        r'<p:cTn\b[^>]*nodeType="clickEffect"[^>]*>.*?</p:cTn>',
                        timing_xml, re.DOTALL,
                    ):
                        block = eff.group(0)
                        has_media = (re.search(r'<p:audio\b', block) is not None
                                     or 'ppaction://media' in block)
                        hit_this_btn = re.search(
                            r'\bspid="' + re.escape(spid) + r'"', block
                        ) is not None
                        if has_media and hit_this_btn:
                            click_effect_bound = True
                            break
            click_to_play = hlink_click_media or click_effect_bound

            # (3) 该 rId 解析到包内 media,且非外部链接
            try:
                with zipfile.ZipFile(file_path) as z:
                    rels_path = 'ppt/slides/_rels/slide6.xml.rels'
                    if audio_rid and rels_path in z.namelist():
                        rc = z.read(rels_path).decode('utf-8')
                        rel_m = re.search(
                            r'<Relationship\b[^>]*Id="' + re.escape(audio_rid)
                            + r'"[^>]*/?>', rc,
                        )
                        if rel_m:
                            rel = rel_m.group(0)
                            is_internal = (
                                'Target="../media/' in rel
                                and 'TargetMode="External"' not in rel
                            )
                            if is_internal:
                                tgt_m = re.search(
                                    r'Target="\.\./(media/[^"]+)"', rel
                                )
                                if tgt_m and ('ppt/' + tgt_m.group(1)) in z.namelist():
                                    audio_packaged = True
            except Exception:
                pass
        if audio_found and click_to_play and audio_packaged:
            dim2[11]["hit"] = True
            dim2[11]["delta"] = 3
            dim2[11]["detail"] = ""
        elif audio_found and not audio_packaged:
            dim2[11]["detail"] = ""
        elif audio_found and not click_to_play:
            dim2[11]["detail"] = ""
        else:
            dim2[11]["detail"] = ""

        # --- 扣分项 ---

        # -5: 第二页PPT整体为一张图片
        # 判定"整体为一张图片": 页面几乎没有可编辑文字(文字被拍平进图片),
        # 且存在覆盖整页的图片对象。办公软件中被转成一张图后正文标题均不可编辑
        s2_editable_text = sum(1 for s in slide2.shapes
                               if s.has_text_frame and any(p.text.strip() for p in s.text_frame.paragraphs))
        s2_fullpage_pic = any(
            s.shape_type == MSO_SHAPE_TYPE.PICTURE
            and s.width and s.height
            and emu_to_cm(s.width) > 25 and emu_to_cm(s.height) > 15
            for s in slide2.shapes
        )
        if s2_editable_text == 0 and s2_fullpage_pic:
            deductions[0]["hit"] = True
            deductions[0]["delta"] = -5
            deductions[0]["detail"] = ""

        # -3: 第三页"教学目标：品格、知识与能力同步生长"及下方文本出现被删除,
        #     或被移动到第四页"重难点与破解路径：把"会做"变成"做得稳""之后
        s3 = prs.slides[2]
        s4 = prs.slides[3]
        s3_texts = get_all_text(s3)
        goal_kw = '教学目标'
        # 第3页是否还保留"教学目标"标题
        goal_in_s3 = any(goal_kw in t for t in s3_texts)
        if not goal_in_s3:
            # 是否被移到第四页(出现在第四页标题"重难点与破解路径"之后)
            s4_shapes_text = [
                ''.join(p.text for p in sp.text_frame.paragraphs)
                for sp in s4.shapes if sp.has_text_frame
            ]
            heavy_idx = next((i for i, t in enumerate(s4_shapes_text) if '重难点与破解路径' in t), None)
            goal_idx = next((i for i, t in enumerate(s4_shapes_text) if goal_kw in t), None)
            moved_after = (heavy_idx is not None and goal_idx is not None and goal_idx > heavy_idx)
            if goal_idx is not None:
                deductions[1]["hit"] = True
                deductions[1]["delta"] = -3
                if moved_after:
                    deductions[1]["detail"] = ""
                else:
                    deductions[1]["detail"] = ""
            else:
                deductions[1]["hit"] = True
                deductions[1]["delta"] = -3
                deductions[1]["detail"] = ""

        # -1: 第3页中出现标题位于页面外
        # 办公软件中"标题位于页面外"指标题框主体落在幻灯片可视区域之外(看不到),
        # 边缘轻微超出不算;以标题与页面矩形的重叠面积占比判定
        s3_title = next(
            (
                shape for shape in s3.shapes
                if '教学目标' in _shape_text(shape)
            ),
            None,
        )
        if (
            s3_title is not None
            and s3_title.left is not None
            and s3_title.top is not None
            and s3_title.width
            and s3_title.height
        ):
            tl, tt = s3_title.left, s3_title.top
            tr, tb = tl + s3_title.width, tt + s3_title.height
            ox = max(0, min(tr, prs.slide_width) - max(tl, 0))
            oy = max(0, min(tb, prs.slide_height) - max(tt, 0))
            visible_area = ox * oy
            title_area = s3_title.width * s3_title.height
            # 标题主体(过半面积)不在页面内,即视为位于页面外
            if title_area > 0 and visible_area < title_area * 0.5:
                deductions[2]["hit"] = True
                deductions[2]["delta"] = -1
                deductions[2]["detail"] = ""

        # -1: 第三页中出现"品格目标"与"知识目标"重叠
        # 按文本定位这两个标题框,判定其矩形是否实质重叠(办公软件中两框相交即视觉重叠)
        def find_shape_by_text(slide, keyword):
            for s in slide.shapes:
                if s.has_text_frame and s.left is not None and s.top is not None and s.width and s.height:
                    if any(keyword in p.text for p in s.text_frame.paragraphs):
                        return s
            return None
        pg = find_shape_by_text(s3, '品格目标')
        kn = find_shape_by_text(s3, '知识目标')
        if pg is not None and kn is not None:
            al, at, ar, ab_ = pg.left, pg.top, pg.left+pg.width, pg.top+pg.height
            bl, bt, br, bb = kn.left, kn.top, kn.left+kn.width, kn.top+kn.height
            ox = max(0, min(ar, br) - max(al, bl))
            oy = max(0, min(ab_, bb) - max(at, bt))
            if ox > 0 and oy > 0:
                deductions[3]["hit"] = True
                deductions[3]["delta"] = -1
                deductions[3]["detail"] = ""

        # -1: 除封面页外出现任意页页码丢失
        # 封面(第1页)不计;第2-6页每页均应显示与其页序对应的页码。
        # 办公软件中页码以文本形式呈现在页面(此处页脚形如"...  ·  N"),
        # 逐页检测是否存在等于该页页序的页码数字
        pg_miss = False
        miss_page = None
        for i in range(1, len(prs.slides)):
            sl = prs.slides[i]
            page_no = str(i + 1)
            found = False
            for s in sl.shapes:
                if not s.has_text_frame:
                    continue
                for p in s.text_frame.paragraphs:
                    t = p.text.strip()
                    if not t:
                        continue
                    # 页脚页码: 文本以该页序数字结尾(如"...· 3"),或整段恰为该页序
                    trailing = re.search(r'(\d+)\s*$', t)
                    if t == page_no or (trailing and trailing.group(1) == page_no):
                        found = True; break
                if found:
                    break
            if not found:
                pg_miss = True; miss_page = i + 1; break
        if pg_miss:
            deductions[4]["hit"] = True
            deductions[4]["delta"] = -1
            deductions[4]["detail"] = ""

        # -1: 第1页或第2页背景为0%的透明度
        # 0%透明度即背景完全不透明(未设置任何透明度),办公软件中表现为背景实色压满、
        # 遮住下层。第1页与第2页任一背景为0%透明度即扣分
        bg_zero_alpha = False
        zero_pages = []
        if slide1_bg:
            a1 = get_blip_alpha(slide1_bg)
            if a1 is None or a1 == 0:
                bg_zero_alpha = True; zero_pages.append('第1页')
        if slide2_bg:
            a2 = get_shape_fill_alpha(slide2_bg)
            if a2 is None or a2 == 0:
                bg_zero_alpha = True; zero_pages.append('第2页')
        if bg_zero_alpha:
            deductions[5]["hit"] = True
            deductions[5]["delta"] = -1
            deductions[5]["detail"] = ""

        # -1: 第1页标题"幼儿游戏活动 课程设计与实施"没有位于页面左侧
        # 按文本定位该标题框,判定其是否位于页面左侧(标题水平中心落在页面左半区)
        s1_title = None
        for s in slide1.shapes:
            if s.has_text_frame and s.left is not None and s.width:
                if any('幼儿游戏活动' in p.text for p in s.text_frame.paragraphs):
                    s1_title = s; break
        if s1_title is not None:
            if s1_title.left + s1_title.width // 2 > prs.slide_width // 2:
                deductions[6]["hit"] = True
                deductions[6]["delta"] = -1
                deductions[6]["detail"] = ""

        # 汇总维度二得分
        result["dim2_items"] = dim2 + deductions
        total = sum(item["delta"] for item in result["dim2_items"])
        result["total_score"] = total
        return result

    except Exception as e:
        result["status"] = "error"
        result["error"] = f"{type(e).__name__}: {e}"
        return result


if __name__ == '__main__':
    # 本地调试入口: 仅用于脚本作者自测,批量运行走 evaluate() 返回值
    # 统一约定: evaluate(dir_path) 接收"脚本所在目录",脚本自行在其中定位待评估文档
    dir_path = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(dir_path), ensure_ascii=False, indent=2))
