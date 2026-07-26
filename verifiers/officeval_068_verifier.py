# -*- coding: utf-8 -*-
# pyright: reportUnknownParameterType=false, reportUnknownVariableType=false, reportUnknownMemberType=false, reportUnknownArgumentType=false, reportMissingParameterType=false
"""
对 "合并_浅蓝可编辑版.pptx" 的自动评估脚本。

评分逻辑：
  维度1（可用与可修改性）：硬性门槛。任一关键项失败 -> 总分 0，且不再检查维度2。
  维度2（完成度评分细则）：
      四个加分项，每项 +5（每项内部所有点都必须满足才计分）
      最终输出命中明细 + 总分。
"""

import json
import os
import re
import sys
import zipfile
from collections import Counter

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 脚本编号（用于返回结果 id 字段）
SCRIPT_ID = "068"

# ---------- 期望的 25 页标题（严格按细则文字）----------
# 顺序：第 1–25 页对应的标题，必须每一页都"命中"对应的名称。
# 注：细则中第 1 页只写明是"封面"，第 8 页只写明是"典型设备组成"，
#    其余页都给出了明确标题，这里完全按细则给出的字面命名。
EXPECTED_TITLES = [
    "封面",                                # 1
    "目录",                                # 2
    "网络定位与构成",                      # 3
    "典型保障场景",                        # 4
    "核心能力价值（一）",                  # 5
    "核心能力价值（二）",                  # 6
    "台站体系与协作架构",                  # 7
    "典型设备组成",                        # 8
    "云岚调度坐席平台",                    # 9
    "内容结构",                            # 10
    "系统总览：调度坐席的四层协作",        # 11
    "登录：进入值守环境",                  # 12
    "主界面：把高频操作放在同一屏",        # 13
    "焦点规则：先确认“操作”，再确认“语音”",# 14
    "工单窗口：先记录，再执行呼叫动作",    # 15
    "线路状态：用可视化面板降低误操作",    # 16
    "接续流程：从来话到三方沟通",          # 17
    "转接与交叉操作：语音不一定跟着窗口走",# 18
    "合并与拆分：处理复杂通话关系",        # 19
    "查号：用更少输入找到正确联系人",      # 20
    "即时会议：临时协同但不污染会议列表",  # 21
    "座席互联：用于班组内快速沟通",        # 22
    "主菜单：按业务动作重新归类",          # 23
    "运行流程与管理要点",                  # 24
    "谢谢聆听",                            # 25
]

# ---------- 颜色相关工具 ----------
def hex_to_rgb(h):
    h = h.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)

def is_light_blue(rgb_hex):
    """浅蓝/浅蓝渐变：R<=G<=B 且 B 偏高、整体偏亮；或落在常见浅蓝色域。"""
    if not rgb_hex:
        return False
    try:
        r, g, b = hex_to_rgb(rgb_hex)
    except Exception:
        return False
    # 整体亮度高（不是深色）
    bright = (r + g + b) / 3 >= 170
    # 蓝色分量占优
    blue_dom = b >= r and b >= g and b - min(r, g) >= 5
    return bright and blue_dom

def is_dark_blue_or_teal(rgb_hex):
    """深蓝 / 蓝绿（用于标题文字判定）。"""
    if not rgb_hex:
        return False
    try:
        r, g, b = hex_to_rgb(rgb_hex)
    except Exception:
        return False
    if (r + g + b) / 3 > 200:        # 太亮的不算
        return False
    # 深蓝：B 显著大于 R
    deep_blue = b >= 90 and b - r >= 25
    # 蓝绿：G、B 都不低且明显大于 R
    teal = g >= 90 and b >= 90 and (g - r) >= 20 and (b - r) >= 20
    return deep_blue or teal

def is_high_saturation_bad(rgb_hex):
    """大面积深黑 / 深红 / 高饱和（与浅蓝风格割裂）。"""
    if not rgb_hex:
        return False
    try:
        r, g, b = hex_to_rgb(rgb_hex)
    except Exception:
        return False
    avg = (r + g + b) / 3
    if avg <= 40:                          # 深黑
        return True
    if r >= 150 and r - g >= 60 and r - b >= 60:   # 深红
        return True
    return False

def shape_fill_hex(sh):
    try:
        if sh.fill.type == 1:              # SOLID
            rgb = sh.fill.fore_color.rgb
            if rgb is not None:
                return str(rgb)
    except Exception:
        pass
    return None

# ---------- 基础读取 ----------
def load_slide_data(prs):
    """采集每页文本、形状、颜色信息。"""
    data = []
    for idx, slide in enumerate(prs.slides, 1):
        shapes_info = []
        all_text_segments = []
        for sh in slide.shapes:
            text = ""
            font_names = []
            font_sizes = []
            bolds = []
            colors = []
            aligns = []
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    if para.alignment is not None:
                        aligns.append(str(para.alignment))
                    for run in para.runs:
                        if run.text:
                            text += run.text
                            font_names.append(run.font.name)
                            font_sizes.append(run.font.size.pt if run.font.size else None)
                            bolds.append(bool(run.font.bold))
                            try:
                                col = run.font.color.rgb
                                colors.append(str(col) if col is not None else None)
                            except Exception:
                                colors.append(None)
                    text += "\n"
                text = text.strip()
            fill_hex = shape_fill_hex(sh)

            try:
                left = sh.left or 0
                top = sh.top or 0
                width = sh.width or 0
                height = sh.height or 0
            except Exception:
                left = top = width = height = 0

            shapes_info.append({
                "shape": sh,
                "type": sh.shape_type,
                "text": text,
                "font_names": font_names,
                "font_sizes": font_sizes,
                "bolds": bolds,
                "text_colors": colors,
                "alignments": aligns,
                "fill_hex": fill_hex,
                "left": left,
                "top": top,
                "width": width,
                "height": height,
            })
            if text:
                all_text_segments.append(text)

        # 背景色
        bg_hex = None
        try:
            if slide.background.fill.type == 1:
                rgb = slide.background.fill.fore_color.rgb
                if rgb is not None:
                    bg_hex = str(rgb)
        except Exception:
            pass

        data.append({
            "index": idx,
            "shapes": shapes_info,
            "texts": all_text_segments,
            "bg_hex": bg_hex,
            "slide": slide,
        })
    return data


# ============================================================
# 维度 1
# ============================================================
def _check_file_openable(file_path):
    """检测文件是否"能够正常打开"。
    返回 (openable: bool, reasons: list[str])。
    检查项：
      (a) 文件存在且非空；
      (b) 是有效的 OOXML 压缩包（zipfile 可解压）；
      (c) 关键部件齐全：至少包含 `ppt/presentation.xml`。
    """
    findings = []

    # (a) 文件存在且非空
    try:
        if not os.path.exists(file_path):
            findings.append("文件不存在")
            return False, findings
        if os.path.getsize(file_path) <= 0:
            findings.append("文件大小为 0")
            return False, findings
    except Exception as e:
        findings.append(f"文件状态读取失败: {e}")
        return False, findings

    # (b)(c) OOXML 结构
    try:
        with zipfile.ZipFile(file_path) as z:
            names = set(z.namelist())
            if "ppt/presentation.xml" not in names:
                findings.append("OOXML 包缺少 ppt/presentation.xml，无法正常打开")
    except zipfile.BadZipFile:
        findings.append("文件不是有效的 OOXML 压缩包（无法作为 .pptx 打开）")
    except Exception as e:
        findings.append(f"打开 .pptx 失败: {e}")

    return len(findings) == 0, findings


def check_dimension1(_prs, _slide_data, file_path):
    """返回 (passed: bool, reasons: list[str])。
    维度 1 现只做"格式 + 能否正常打开"的硬门槛，因此 prs / slide_data 在此不使用，
    仅为与主流程调用签名保持一致而保留。"""
    reasons = []

    # 1. 格式：交付文件为 .pptx
    if not file_path.lower().endswith(".pptx"):
        reasons.append("文件扩展名不是 .pptx")
        return False, reasons

    # 2. 能否正常打开
    openable, findings = _check_file_openable(file_path)
    if not openable:
        reasons.append("文件无法正常打开：" + "；".join(findings))
        return False, reasons

    reasons.append("文件为 .pptx，能够正常打开")
    return True, reasons


# ============================================================
# 维度 2
# ============================================================
def slide_title_text(info):
    """取一页中最可能是"标题"的文本：选纵向靠上、字号较大、且不是 '01'/'页码' 之类的短编号。"""
    candidates = []
    for sh in info["shapes"]:
        t = sh["text"].strip()
        if not t:
            continue
        if re.fullmatch(r"\d{1,3}", t):           # 纯数字（页码）
            continue
        if re.fullmatch(r"PDF原第\d+页|PPT原第\d+页", t):  # 旧标记（如果存在）
            continue
        max_size = max([s for s in sh["font_sizes"] if s] or [0])
        candidates.append((sh["top"], -max_size, t, sh))
    candidates.sort()
    # 取前 6 个合并作为标题搜索池
    return " ".join(c[2] for c in candidates[:6])


def _norm(s):
    """归一化字符串：去空白、统一中英文标点，方便做包含匹配。"""
    if not s:
        return ""
    s = s.replace("　", "").replace(" ", "").replace("\n", "").replace("\t", "")
    # 引号 / 冒号 / 括号 全部转半角
    table = str.maketrans({
        "“": '"', "”": '"', "‘": "'", "’": "'",
        "：": ":", "（": "(", "）": ")",
        "，": ",", "。": ".", "、": ",",
    })
    return s.translate(table)


def check_plus1_pages_and_order(prs, slide_data, file_path):
    """+5：交付 .pptx，共 25 页；每一页按细则给出的顺序与标题一一对应。"""
    detail = []

    # (a) 文件格式
    if not file_path.lower().endswith(".pptx"):
        return False, ["文件不是 .pptx"]

    # (b) 25 页
    if len(prs.slides) != 25:
        return False, [f"页数 {len(prs.slides)} ≠ 25"]

    # (c) 每页标题必须匹配细则
    #     细则写的是"第 1 页 = 封面、第 2 页 = 目录、第 10 页 = 内容结构"——
    #     其中第 1 页是"封面"角色，第 2 页必须是"目录"（不能与第 10 页"内容结构"混用）。
    #     其余页直接按字面匹配。
    mismatches = []
    for i, expected in enumerate(EXPECTED_TITLES, 1):
        info = slide_data[i - 1]
        page_text_norm = _norm("\n".join(info["texts"]))

        if i == 1:                            # 封面
            # 不允许出现"目录 / CONTENTS / 内容结构"等目录性关键词
            if ("目录" in page_text_norm
                    or "contents" in page_text_norm.lower()
                    or "内容结构" in page_text_norm):
                mismatches.append((i, expected, "含目录关键字，非封面"))
                continue

            # 必须能验证"封面身份"：满足以下任一：
            #   1) 版式名 / 母版名 包含 "封面" / "title" / "cover"
            #   2) 存在 PP_PLACEHOLDER.TITLE / CENTER_TITLE / SUBTITLE 版式占位符
            #   3) 存在文字命中封面关键标题（"云岚" / "调度坐席" / "平台" 等主标题词），
            #      且有明显的大字号主标题（>=24pt）
            slide = info["slide"]
            cover_id_ok = False
            reasons_cover = []

            # (1) 版式 / 母版名
            try:
                layout_name = (slide.slide_layout.name or "").lower()
            except Exception:
                layout_name = ""
            try:
                master_name = (slide.slide_layout.slide_master.name or "").lower()
            except Exception:
                master_name = ""
            if any(k in layout_name for k in ("封面", "cover", "title"))\
                    or any(k in master_name for k in ("封面", "cover", "title")):
                cover_id_ok = True
                reasons_cover.append(f"layout={layout_name!r}")

            # (2) 标题/副标题占位符
            if not cover_id_ok:
                try:
                    from pptx.enum.shapes import PP_PLACEHOLDER
                    ph_kinds = set()
                    for shp in slide.placeholders:
                        try:
                            ph_kinds.add(shp.placeholder_format.type)
                        except Exception:
                            pass
                    if (PP_PLACEHOLDER.TITLE in ph_kinds
                            or PP_PLACEHOLDER.CENTER_TITLE in ph_kinds
                            or PP_PLACEHOLDER.SUBTITLE in ph_kinds):
                        cover_id_ok = True
                        reasons_cover.append("含 TITLE/CENTER_TITLE/SUBTITLE 占位符")
                except Exception:
                    pass

            # (3) 关键标题词 + 大字号
            if not cover_id_ok:
                cover_keywords = ("云岚", "调度坐席", "调度台", "坐席平台")
                if any(k in page_text_norm for k in cover_keywords):
                    max_size = 0
                    for sh in info["shapes"]:
                        sizes = [s for s in sh["font_sizes"] if s] or [0]
                        if max(sizes) > max_size:
                            max_size = max(sizes)
                    if max_size >= 24:
                        cover_id_ok = True
                        reasons_cover.append(f"含封面关键词且主标题字号={max_size}")

            if not cover_id_ok:
                mismatches.append((i, expected, "未能识别为封面（版式/占位符/关键标题均不匹配）"))
            continue

        if i == 2:                            # 目录
            # 严格匹配"目录 / CONTENTS"，不接受"内容结构"（那是第 10 页）
            if not any(k in page_text_norm for k in ("目录",))\
                    and "contents" not in page_text_norm.lower():
                mismatches.append((i, expected, "未出现'目录'/'CONTENTS'"))
            continue

        # 其余页：字面包含匹配
        exp_norm = _norm(expected)
        if exp_norm not in page_text_norm:
            mismatches.append((i, expected, "字面标题未命中"))

    if mismatches:
        for item in mismatches:
            if len(item) == 3:
                i, t, why = item
                detail.append(f"第 {i} 页未匹配到标题 {t!r}：{why}")
            else:
                i, t = item
                detail.append(f"第 {i} 页未匹配到标题: {t!r}")
        return False, detail

    detail.append("文件为 .pptx；共 25 页；25 页标题与顺序全部匹配细则（第 2 页严格为目录，第 1 页身份为封面）")
    return True, detail


def check_plus2_editable(prs, slide_data):
    """+5：PDF 转换页面文字 + 全部 25 页可编辑。
    细则要求：
      - 第 1–8 页 及 第 24、25 页：
          文本类：标题、正文、编号、说明 —— 均为可编辑文本框；
          形状类：流程框、箭头、圆形、图标、分隔线、示意结构 —— 均为可编辑形状或
                 可编辑组合对象。
      - 第 9–23 页：文字、表格、形状、流程对象、图标 保持可编辑。

    本函数**按对象类别**逐类检测，不再仅用数量阈值近似：
      * "标题"     = 最大字号非纯数字文本；
      * "正文"     = 中/大量文字（≥ 8 个字符）、字号不属于"标题字号"的文本框；
      * "编号"     = 短小的编号型文本（纯数字、"01"、"①-⑳"、"Step X"、"1." 等）；
      * "说明"     = 明显小字号（<= 14pt 或页面最小字号档）的补充文本；
      * "流程框"   = 矩形 / 圆角矩形 类 AUTO_SHAPE；
      * "箭头"     = ARROW / CHEVRON / PENTAGON 等箭头族 AUTO_SHAPE，
                    或带箭头端点的 LINE / CONNECTOR；
      * "圆形"     = OVAL / ELLIPSE；
      * "图标"     = 小尺寸 AUTO_SHAPE / FREEFORM / GROUP（长宽 ≤ 页面 1/6），
                    或非整页的矢量小图片；
      * "分隔线"   = LINE / CONNECTOR，或极细长的矩形；
      * "示意结构" = GROUP，或"多个可编辑形状 + 线/箭头"组合。
    """
    detail = []
    PDF_PAGES = set(list(range(1, 9)) + [24, 25])     # 1-8、24、25
    PPT_PAGES = set(range(9, 24))                     # 9-23

    # 兼容不同 python-pptx 版本
    def _mst(name):
        return getattr(MSO_SHAPE_TYPE, name, None)

    T_PICTURE   = _mst("PICTURE")
    T_AUTO      = _mst("AUTO_SHAPE")
    T_FREEFORM  = _mst("FREEFORM")
    T_GROUP     = _mst("GROUP")
    T_LINE      = _mst("LINE")

    def _auto_shape_name(sh_dict):
        """返回该 shape 的 auto_shape_type 名称（大写字符串），失败时 ''。"""
        try:
            ast = sh_dict["shape"].auto_shape_type
            return str(ast).upper() if ast is not None else ""
        except Exception:
            return ""

    def is_full_page_picture(s):
        return (s["type"] == T_PICTURE
                and s["width"] >= prs.slide_width * 0.85
                and s["height"] >= prs.slide_height * 0.85)

    # ---------- 文本类分类 ----------
    NUMBERING_RE = re.compile(
        (
            r"^(?:"
            r"\d{1,3}"                                # 1 / 01 / 100
            r"|[①-⑳]"                                # 圈号
            r"|[一二三四五六七八九十]{1,3}"           # 中文序号
            r"|(?:step|步骤|序号)\s*\d{1,3}"          # Step 1 / 步骤1
            r"|\d{1,3}[.、)）]"                       # 1. / 1) / 1、
            r"|[A-Za-z][.、)）]"                      # A. / a)
            r")$"
        ),
        re.IGNORECASE,
    )

    def classify_text_shape(sh_dict, title_size):
        """把一个含文字的 shape 归入 title / body / number / caption / other。"""
        text = sh_dict["text"].strip()
        if not text:
            return None
        sizes = [s for s in sh_dict["font_sizes"] if s]
        max_size = max(sizes) if sizes else 0
        text_wo_ws = re.sub(r"\s+", "", text)

        # 编号：短文本且匹配编号正则
        if len(text_wo_ws) <= 6 and NUMBERING_RE.match(text_wo_ws):
            return "number"
        # 标题：字号 == 页面最大字号，且不是编号
        if title_size and max_size and max_size >= title_size - 0.5:
            return "title"
        # 说明：小字号（<=14pt）
        if max_size and max_size <= 14:
            return "caption"
        # 正文：其余较长文字
        if len(text_wo_ws) >= 8:
            return "body"
        return "other"

    # ---------- 形状类分类 ----------
    def classify_non_text_shape(sh_dict):
        """返回该 shape 所属的形状类别集合（可能多类），非可编辑形状返回空集。"""
        t = sh_dict["type"]
        if t == T_PICTURE:
            # 图片默认不算"可编辑形状"，但小矢量图片可作为"图标"的兜底
            if (sh_dict["width"] <= prs.slide_width / 6
                    and sh_dict["height"] <= prs.slide_height / 6):
                return {"icon_pic"}
            return set()

        cats = set()
        name = _auto_shape_name(sh_dict)
        w, h = sh_dict["width"], sh_dict["height"]
        page_w, page_h = prs.slide_width, prs.slide_height

        # 分组：本身算"示意结构"
        if t == T_GROUP:
            cats.add("diagram")

        # 直线 / 连接线
        if t == T_LINE or "CONNECTOR" in name or "LINE" in name:
            cats.add("divider")
            # 若形状名中含 ARROW，也当作箭头
            if "ARROW" in name:
                cats.add("arrow")

        # AUTO_SHAPE 细分
        if t == T_AUTO or t == T_FREEFORM:
            if "ARROW" in name or "CHEVRON" in name or "PENTAGON" in name:
                cats.add("arrow")
            if "OVAL" in name or "ELLIPSE" in name or "CIRCLE" in name:
                cats.add("circle")
            if "RECTANGLE" in name or "ROUND" in name or name == "RECT":
                # 极细长的矩形当分隔线
                if (w > 0 and h > 0
                        and (min(w, h) < max(page_w, page_h) * 0.01
                             or (max(w, h) / max(min(w, h), 1)) > 25)):
                    cats.add("divider")
                else:
                    cats.add("flowbox")
            # 小尺寸形状 → 图标
            if (w > 0 and h > 0
                    and w <= page_w / 6 and h <= page_h / 6):
                cats.add("icon")
            # FREEFORM 通常是图标 / 示意
            if t == T_FREEFORM:
                cats.add("icon")

        return cats

    failures = []

    for info in slide_data:
        idx = info["index"]
        shapes = info["shapes"]

        # 是否含整页图片（PDF 截图特征）
        has_full_page_pic = any(is_full_page_picture(s) for s in shapes)
        if has_full_page_pic:
            failures.append(f"第{idx}页存在整页图片（PDF 截图特征，未做可编辑转换）")
            continue

        # 计算该页最大字号（用于识别"标题"）
        page_max_size = 0
        for s in shapes:
            if not s["text"]:
                continue
            if re.fullmatch(r"\d{1,3}", s["text"].strip()):
                continue
            sizes = [z for z in s["font_sizes"] if z]
            if sizes and max(sizes) > page_max_size:
                page_max_size = max(sizes)

        # —— 文本类逐类归档 ——
        text_cats = {"title": [], "body": [], "number": [], "caption": [], "other": []}
        for s in shapes:
            if not s["text"]:
                continue
            k = classify_text_shape(s, page_max_size)
            if k:
                text_cats[k].append(s)

        # —— 形状类逐类归档 ——
        shape_cats = {"flowbox": [], "arrow": [], "circle": [],
                      "icon": [], "divider": [], "diagram": []}
        for s in shapes:
            if s["text"] and s["type"] not in (T_GROUP,):
                # 已归入文本类的形状不再重复计入"流程框/箭头/圆形"等
                # 但带文字的组合仍可算"示意结构"
                continue
            cats = classify_non_text_shape(s)
            for c in cats:
                if c == "icon_pic":
                    shape_cats["icon"].append(s)
                elif c in shape_cats:
                    shape_cats[c].append(s)
        # 若页内含 GROUP 或"多个 flowbox+箭头/线"组合，则视为存在示意结构
        if (not shape_cats["diagram"]
                and len(shape_cats["flowbox"]) >= 2
                and (shape_cats["arrow"] or shape_cats["divider"])):
            shape_cats["diagram"] = ["<derived-from-flowbox+line>"]

        # ================= 判定 =================
        if idx in PDF_PAGES:
            missing_text = [n for n in ("title", "body", "number", "caption")
                            if not text_cats[n]]
            missing_shape = [n for n in ("flowbox", "arrow", "circle",
                                          "icon", "divider", "diagram")
                             if not shape_cats[n]]
            if missing_text:
                cn = {"title": "标题", "body": "正文",
                      "number": "编号", "caption": "说明"}
                failures.append(
                    f"第{idx}页缺少可编辑文本类别: "
                    + "、".join(cn[k] for k in missing_text)
                )
                continue
            if missing_shape:
                cn = {"flowbox": "流程框", "arrow": "箭头", "circle": "圆形",
                      "icon": "图标", "divider": "分隔线", "diagram": "示意结构"}
                failures.append(
                    f"第{idx}页缺少可编辑形状类别: "
                    + "、".join(cn[k] for k in missing_shape)
                )
                continue

        elif idx in PPT_PAGES:
            # 文字：至少存在"标题 + 正文/说明"其一
            if not text_cats["title"]:
                failures.append(f"第{idx}页缺少可编辑标题文字")
                continue
            if not (text_cats["body"] or text_cats["caption"] or text_cats["number"]):
                failures.append(f"第{idx}页除标题外无其他可编辑文字")
                continue
            # 形状 / 流程对象 / 图标：至少存在两类
            present = [k for k, v in shape_cats.items() if v]
            if len(present) < 2:
                failures.append(
                    f"第{idx}页可编辑形状类别过少（仅 {present}），需同时具备形状/流程对象/图标"
                )
                continue
            # 表格：若页内有大幅图片（可能是被烧录的表格/示意），判失败
            for s in shapes:
                if (s["type"] == T_PICTURE
                        and s["width"] > prs.slide_width * 0.4
                        and s["height"] > prs.slide_height * 0.4):
                    failures.append(
                        f"第{idx}页存在大幅图片（{int(s['width'])}×{int(s['height'])}），可能为不可编辑表格/示意"
                    )
                    break
            else:
                # 表格：如果 python-pptx 识别到 TABLE，则要求它是真表格（可编辑）
                #       如果没有 TABLE，也不强制（细则语义："如有表格则必须可编辑"）
                pass

    if failures:
        for f in failures:
            detail.append(f)
        return False, detail

    detail.append(
        "第 1-8、24、25 页：标题/正文/编号/说明 均为可编辑文本框；"
        + "流程框/箭头/圆形/图标/分隔线/示意结构 均为可编辑形状或组合"
    )
    detail.append("第 9-23 页：文字、表格、形状、流程对象、图标 均保持可编辑")
    return True, detail


def check_plus3_style(prs, slide_data):
    """+5：视觉风格细则。逐条对应：
      (a) 全部 25 页背景统一采用浅蓝色或浅蓝渐变背景；
      (b) 不出现与整体风格割裂的深黑、深红或高饱和背景；
      (c) 标题统一使用深蓝色或蓝绿色文字；
      (d) 标题字体、字号、加粗 和 左对齐方式保持一致；
      (e) 正文中中文字体、英文字体、字号、行距和项目符号样式统一；
      (f) 正文清晰可读；
      (g) 使用白色、淡蓝色或蓝绿色圆角框；
      (h) 圆角框的边框颜色、线宽和圆角程度保持一致。
    """
    detail = []
    problems = []

    # ===== (a) + (b) 背景 =====
    # 背景来源优先级：slide 自身 -> slide_layout -> slide_master
    # 不再接受纯白（FFFFFF）作为浅蓝色的替代——rubric 明确要求"浅蓝色或浅蓝渐变"。
    LIGHT_BLUE_WHITELIST = {"F7FBFF", "F2F8FF", "EAF4FF", "E6F2FF", "DDEEFF"}

    def _fill_hex_from_fill(fill):
        """从 python-pptx 的 fill 对象解析出主色（SOLID/GRADIENT 均支持）。返回 (hex, is_gradient)。"""
        try:
            ft = fill.type
        except Exception:
            return None, False
        # SOLID
        try:
            if ft == 1:                       # MSO_FILL.SOLID
                rgb = fill.fore_color.rgb
                if rgb is not None:
                    return str(rgb), False
        except Exception:
            pass
        # GRADIENT —— python-pptx 未提供高层 API，读 XML 里的 <a:gradFill><a:gs> 停靠色
        try:
            xml = fill._xPr.xml if hasattr(fill, "_xPr") else ""
        except Exception:
            xml = ""
        if "<a:gradFill" in xml or ft == 3:  # 3 = GRADIENT (若 python-pptx 暴露)
            m: "list[str]" = re.findall(r'<a:srgbClr\s+val="([0-9A-Fa-f]{6})"', xml)
            if m:
                # 取渐变色中"最浅"的一档做代表色
                def brightness(h: str) -> int:
                    r, g, b = hex_to_rgb(h)
                    return r + g + b
                pick: str = max(m, key=brightness)
                return pick.upper(), True
        return None, False

    def _resolve_slide_bg(slide):
        """按 slide -> layout -> master 顺序取有效背景色。返回 (hex, source, is_gradient)。"""
        for source_name, obj in (
            ("slide", slide),
            ("layout", getattr(slide, "slide_layout", None)),
            ("master", getattr(getattr(slide, "slide_layout", None), "slide_master", None)),
        ):
            if obj is None:
                continue
            try:
                bg_fill = obj.background.fill
            except Exception:
                continue
            hex_val, is_grad = _fill_hex_from_fill(bg_fill)
            if hex_val:
                return hex_val.upper(), source_name, is_grad
        return None, None, False

    bg_bad = []
    bg_harsh = []
    bg_unknown = []
    bg_hex_seen = Counter()
    for info in slide_data:
        slide = info["slide"]
        bg_hex, source, _is_grad = _resolve_slide_bg(slide)
        if bg_hex is None:
            bg_unknown.append(info["index"])
            continue
        bg_hex_seen[bg_hex] += 1
        if is_high_saturation_bad(bg_hex):
            bg_harsh.append((info["index"], bg_hex))
            continue
        if not (is_light_blue(bg_hex) or bg_hex in LIGHT_BLUE_WHITELIST):
            bg_bad.append((info["index"], bg_hex, source))

    if bg_unknown:
        problems.append(f"(a) 存在无法判定背景色的页: {bg_unknown}（要求全部 25 页背景可判定且统一）")
    if bg_bad:
        problems.append(f"(a) 背景非浅蓝/浅蓝渐变（含 FFFFFF 纯白）: {bg_bad}")
    if bg_harsh:
        problems.append(f"(b) 出现深黑/深红/高饱和背景: {bg_harsh}")
    # 统一性：只允许 1 种主背景色，且该主背景色覆盖全部 25 页（渐变的浅色停靠色相同也算统一）
    if bg_hex_seen and len(slide_data) == 25:
        _top_bg, top_cnt = bg_hex_seen.most_common(1)[0]
        if top_cnt < 25:
            problems.append(
                f"(a) 背景色未在全部 25 页保持统一: 分布={dict(bg_hex_seen)}"
            )

    # ===== 选出每页"标题"和"正文" runs =====
    title_runs = []   # 每条: (slide_idx, shape, run_font_name, run_size, run_bold, color, align)
    body_runs  = []

    for info in slide_data:
        # 用最大字号的非数字文本作为该页标题
        biggest_size = 0
        biggest = None
        for sh in info["shapes"]:
            if not sh["text"]:
                continue
            if re.fullmatch(r"\d{1,3}", sh["text"].strip()):
                continue
            sizes = [s for s in sh["font_sizes"] if s] or [0]
            if max(sizes) > biggest_size:
                biggest_size = max(sizes)
                biggest = sh

        for sh in info["shapes"]:
            if not sh["text"]:
                continue
            is_title = (sh is biggest)
            tf = sh["shape"].text_frame
            for p_idx, para in enumerate(tf.paragraphs):
                align = str(para.alignment) if para.alignment is not None else None
                line_spacing = None
                try:
                    line_spacing = para.line_spacing
                except Exception:
                    pass
                # 项目符号检测：段落 xml 中是否带 <a:buChar/<a:buAutoNum/<a:buNone>
                try:
                    pPr_xml = para._pPr.xml if para._pPr is not None else ""
                except Exception:
                    pPr_xml = ""
                bullet_kind = "none"
                if "<a:buChar" in pPr_xml: bullet_kind = "char"
                elif "<a:buAutoNum" in pPr_xml: bullet_kind = "num"
                elif "<a:buBlip" in pPr_xml: bullet_kind = "pic"
                for run in para.runs:
                    if not run.text:
                        continue
                    name = run.font.name
                    size = run.font.size.pt if run.font.size else None
                    bold = bool(run.font.bold)
                    try:
                        col = run.font.color.rgb
                        col = str(col) if col is not None else None
                    except Exception:
                        col = None
                    rec = {
                        "slide": info["index"],
                        "name": name,
                        "size": size,
                        "bold": bold,
                        "color": col,
                        "align": align,
                        "line_spacing": line_spacing,
                        "bullet": bullet_kind,
                        "text": run.text,
                    }
                    if is_title:
                        title_runs.append(rec)
                    else:
                        body_runs.append(rec)

    # ===== (c) 标题颜色：深蓝 / 蓝绿（不再放行 FFFFFF；白色标题只能在深色底色块上出现，
    #        但 rubric 明确要求"深蓝或蓝绿" -> 白色一律判失败） =====
    title_color_bad = []
    title_color_unknown = []
    for r in title_runs:
        col = r["color"]
        if col is None:
            title_color_unknown.append((r["slide"], r["text"][:10]))
            continue
        if not is_dark_blue_or_teal(col):
            title_color_bad.append((r["slide"], col, r["text"][:10]))
    if title_color_bad:
        problems.append(f"(c) 标题颜色非深蓝/蓝绿: {title_color_bad[:10]}")
    if title_color_unknown and len(title_color_unknown) > len(title_runs) * 0.2:
        problems.append(
            f"(c) 过多标题的颜色无法从 run 上读取（占比 {len(title_color_unknown)}/{len(title_runs)}）"
        )

    # ===== (d) 标题字体、字号、加粗、左对齐 一致（阈值收紧到 0.95） =====
    title_fonts  = Counter(r["name"] for r in title_runs if r["name"])
    title_sizes  = Counter(round(r["size"]) for r in title_runs if r["size"])
    title_bolds  = Counter(r["bold"] for r in title_runs)
    title_aligns = Counter(r["align"] for r in title_runs if r["align"] is not None)

    def _top_ratio(c):
        if not c: return 1.0
        return c.most_common(1)[0][1] / sum(c.values())

    # rubric 要求"保持一致"——阈值 0.95，仅允许极少数版式差异（如封面/结尾页可能不同）
    UNIFORM = 0.95

    if title_fonts and _top_ratio(title_fonts) < UNIFORM:
        problems.append(f"(d) 标题字体不统一: {title_fonts.most_common(5)}")
    if title_sizes and _top_ratio(title_sizes) < UNIFORM:
        problems.append(f"(d) 标题字号不统一: {title_sizes.most_common(5)}")
    if title_bolds and _top_ratio(title_bolds) < UNIFORM:
        problems.append(f"(d) 标题加粗不统一: {dict(title_bolds)}")
    # 左对齐：rubric 明确"左对齐"——只接受 LEFT；CENTER/RIGHT/JUSTIFY 均判失败
    if title_aligns:
        non_left = {k: v for k, v in title_aligns.items() if "LEFT" not in (k or "")}
        left_cnt = sum(v for k, v in title_aligns.items() if "LEFT" in (k or ""))
        total = sum(title_aligns.values())
        # 允许极少数封面/结尾等特殊版式的例外（例如"封面"和"谢谢聆听"两页 -> 2/25 ≈ 8%）
        if left_cnt / total < 0.85 or non_left:
            problems.append(f"(d) 标题对齐方式非左对齐或不一致: {dict(title_aligns)}")

    # ===== (e) 正文中文字体 / 英文字体 / 字号 / 行距 / 项目符号 统一 =====
    cn_fonts = Counter()
    en_fonts = Counter()
    body_sizes = Counter(round(r["size"]) for r in body_runs if r["size"])
    body_line_spacings = Counter()
    body_bullets = Counter()
    for r in body_runs:
        # 中/英文字体粗略区分：name 含 CJK 关键字 或 文本里多为中文 -> 视为中文字体
        name = r["name"]
        text = r["text"]
        if name:
            cn_chars = sum(1 for ch in text if "一" <= ch <= "鿿")
            if cn_chars >= max(1, len(text) // 2):
                cn_fonts[name] += 1
            else:
                en_fonts[name] += 1
        if r["line_spacing"] is not None:
            body_line_spacings[r["line_spacing"]] += 1
        body_bullets[r["bullet"]] += 1

    if cn_fonts and _top_ratio(cn_fonts) < UNIFORM:
        problems.append(f"(e) 正文中文字体不统一: {cn_fonts.most_common(5)}")
    if en_fonts and _top_ratio(en_fonts) < UNIFORM:
        problems.append(f"(e) 正文英文字体不统一: {en_fonts.most_common(5)}")
    if body_sizes and _top_ratio(body_sizes) < 0.85:
        # 正文字号允许有次级层次（如"正文 + 说明"两档），阈值略放宽
        problems.append(f"(e) 正文字号不统一: {body_sizes.most_common(5)}")
    if body_line_spacings and _top_ratio(body_line_spacings) < UNIFORM:
        problems.append(f"(e) 正文行距不统一: {body_line_spacings.most_common(5)}")
    if body_bullets and _top_ratio(body_bullets) < UNIFORM:
        problems.append(f"(e) 项目符号样式不统一: {dict(body_bullets)}")

    # ===== (f) 正文清晰可读 =====
    #   判定：正文字号不应过小，正文文字颜色与背景对比度不应过低
    illegible = []
    for r in body_runs:
        if r["size"] and r["size"] < 8:
            illegible.append((r["slide"], "size", r["size"]))
        # 颜色对比：浅色文字 + 浅色背景 -> 不可读
        if r["color"]:
            try:
                rr, gg, bb = hex_to_rgb(r["color"])
                if (rr + gg + bb) / 3 > 220:           # 接近白色
                    bg = slide_data[r["slide"] - 1]["bg_hex"]
                    if bg:
                        br, bgc, bbc = hex_to_rgb(bg)
                        if (br + bgc + bbc) / 3 > 220:
                            illegible.append((r["slide"], "contrast", r["text"][:8]))
            except Exception:
                pass
    if len(illegible) > 5:
        problems.append(f"(f) 正文可读性存在问题（{len(illegible)} 处）: {illegible[:5]}")

    # ===== (g) (h) 圆角框：颜色（白/淡蓝/蓝绿）+ 边框颜色 / 线宽 / 圆角程度一致 =====
    rounded_fills = Counter()
    rounded_line_colors = Counter()
    rounded_line_widths = Counter()
    rounded_adjustments = Counter()           # 圆角程度
    bad_rounded_fill = []
    for info in slide_data:
        for sh in info["shapes"]:
            try:
                ast = sh["shape"].auto_shape_type
            except Exception:
                ast = None
            if ast is None or "ROUND" not in str(ast):
                continue
            # 填充色
            fill_hex = sh["fill_hex"]
            if fill_hex:
                f_up = fill_hex.upper()
                rounded_fills[f_up] += 1
                ok_color = (f_up == "FFFFFF") or is_light_blue(f_up) or is_dark_blue_or_teal(f_up)
                if not ok_color:
                    bad_rounded_fill.append((info["index"], f_up))
            # 边框颜色 / 线宽
            try:
                line = sh["shape"].line
                try:
                    lc = line.color.rgb
                    if lc is not None:
                        rounded_line_colors[str(lc).upper()] += 1
                except Exception:
                    pass
                try:
                    if line.width is not None:
                        rounded_line_widths[int(line.width)] += 1
                except Exception:
                    pass
            except Exception:
                pass
            # 圆角程度：从 <a:prstGeom> 下的 <a:avLst><a:gd .. fmla="val NNN"/> 取值
            try:
                xml = sh["shape"]._element.xml
                m = re.search(r'<a:gd[^/]*fmla="val (\d+)"', xml)
                if m:
                    rounded_adjustments[int(m.group(1))] += 1
            except Exception:
                pass

    if bad_rounded_fill:
        problems.append(f"(g) 圆角框填充色不在白/淡蓝/蓝绿: {bad_rounded_fill[:10]}")
    if rounded_line_colors and _top_ratio(rounded_line_colors) < UNIFORM:
        problems.append(f"(h) 圆角框边框颜色不一致: {rounded_line_colors.most_common(5)}")
    if rounded_line_widths and _top_ratio(rounded_line_widths) < UNIFORM:
        problems.append(f"(h) 圆角框线宽不一致: {rounded_line_widths.most_common(5)}")
    if rounded_adjustments and _top_ratio(rounded_adjustments) < 0.85:
        problems.append(f"(h) 圆角程度不一致: {rounded_adjustments.most_common(5)}")

    if problems:
        for p in problems:
            detail.append(p)
        return False, detail

    detail.append("(a)(b) 全部 25 页背景为浅蓝/浅蓝渐变，无深黑/深红/高饱和背景")
    detail.append("(c)(d) 标题统一为深蓝/蓝绿，字体/字号/加粗/对齐方式一致")
    detail.append("(e)(f) 正文中英文字体、字号、行距、项目符号样式统一，正文清晰可读")
    detail.append("(g)(h) 圆角框颜色为白/淡蓝/蓝绿，边框颜色、线宽与圆角程度一致")
    return True, detail


def check_plus4_page_numbers(prs, slide_data):
    """+5：全部 25 页页码：按照最终合并顺序连续显示 1 至 25。
    逐点对应：
      (a) 全部 25 页都要显示页码；
      (b) 页码按合并后顺序，依次连续显示 1, 2, ..., 25；
      (c) 同一页只应显示一个页码（若同页出现多个不同页码值，视为不连续）。
    判定"页码型文本框"的口径：
      - 文本框只包含一行；
      - 内容为纯数字 / '01' 这种 0 前导两位 / '第X页' / 'Page X' / 'X/25'；
      - 位于页面边缘（左右靠边或上下靠边的 15% 区域）—— 用于把目录里的章节编号排除掉。
    """
    detail = []
    problems = []

    sw, sh_ = prs.slide_width, prs.slide_height
    EDGE = 0.18                          # 距离页面边缘 18% 视为"边角"

    def parse_page_number_text(t):
        """如果该字符串看起来像一枚'页码'，返回其数字，否则返回 None。"""
        t = t.strip()
        if not t:
            return None
        if re.fullmatch(r"0?\d{1,3}", t):
            return int(t)
        m = re.fullmatch(r"第\s*(\d{1,3})\s*页", t)
        if m:
            return int(m.group(1))
        m = re.fullmatch(r"(?:page|p\.?)\s*(\d{1,3})", t, re.IGNORECASE)
        if m:
            return int(m.group(1))
        m = re.fullmatch(r"(\d{1,3})\s*/\s*(\d{1,3})", t)
        if m:
            return int(m.group(1))
        return None

    def shape_is_at_edge(sh):
        cx = sh["left"] + sh["width"] / 2
        cy = sh["top"]  + sh["height"] / 2
        return (cx < sw * EDGE or cx > sw * (1 - EDGE)
                or cy < sh_ * EDGE or cy > sh_ * (1 - EDGE))

    def collect_page_number_shapes(info):
        """返回该页所有页码形状的 (数字值, shape 描述)。"""
        hits = []
        for sh in info["shapes"]:
            txt = sh["text"].strip()
            if not txt:
                continue
            lines = [ln.strip() for ln in txt.splitlines() if ln.strip()]
            if len(lines) != 1:
                continue
            n = parse_page_number_text(lines[0])
            if n is None:
                continue
            # 必须位于页面边缘，否则视作目录/正文里的编号
            if not shape_is_at_edge(sh):
                continue
            hits.append((n, sh))
        return hits

    # —— (a)(b)(c) 逐页检查页码 ——
    missing_pages = []
    wrong_pages   = []
    duplicates_on_page = []

    for info in slide_data:
        idx = info["index"]
        page_num_shapes = collect_page_number_shapes(info)
        nums = [n for n, _ in page_num_shapes]

        if not page_num_shapes:
            missing_pages.append(idx)
            continue

        # (b) 页码值必须等于当前页索引
        if idx not in nums:
            wrong_pages.append((idx, nums))

        # (c) 同一页出现多个不同页码值（视为不连续）
        unique_vals = set(nums)
        if len(unique_vals) > 1:
            duplicates_on_page.append((idx, sorted(unique_vals)))

    if missing_pages:
        problems.append(f"(a) 以下页未显示页码: {missing_pages}")
    if wrong_pages:
        problems.append(f"(b) 以下页页码与最终顺序不符: {wrong_pages[:10]}")
    if duplicates_on_page:
        problems.append(f"(c) 同页出现多个不同页码值（非连续单一）: {duplicates_on_page[:10]}")

    if problems:
        for p in problems:
            detail.append(p)
        return False, detail

    detail.append("全部 25 页按合并顺序连续显示 1..25")
    return True, detail


# ============================================================
# 主流程
# ============================================================
# 被评估的目标文档名（脚本自己在传入的目录里定位并打开）
TARGET_FILE_NAME = "合并_浅蓝可编辑版.pptx"


def evaluate(dir_path: str) -> "dict[str, object]":
    """统一入口：接收"脚本所在目录的路径"，脚本自己在该目录里定位并打开被评估文档。

    返回结构化字典（见 §2.2）。
    """
    result: "dict[str, object]" = {
        "id": SCRIPT_ID,
        "file_name": TARGET_FILE_NAME,
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": 20,
    }

    try:
        file_path = os.path.join(dir_path, TARGET_FILE_NAME)
        if not os.path.exists(file_path):
            result["status"] = "error"
            result["error"] = f"文件不存在: {file_path}"
            return result

        prs = Presentation(file_path)
        slide_data = load_slide_data(prs)

        # ------- 维度 1 -------
        d1_pass, d1_reasons = check_dimension1(prs, slide_data, file_path)
        result["dim1_pass"] = d1_pass
        if not d1_pass:
            result["dim1_reason"] = "；".join(d1_reasons)
            result["total_score"] = 0
            return result

        # ------- 维度 2 -------
        plus_items = [
            (
                5,
                "最终文件格式与页数：交付.pptx文件，共25页。第1至第8页，顺序为封面、目录、网络定位与构成、典型保障场景、核心能力价值（一）、核心能力价值（二）、台站体系与协作架构、典型设备组成；第9至第20页依次对应“云岚调度坐席平台”、“内容结构”、“系统总览：调度坐席的四层协作”、“登录：进入值守环境”、“主界面：把高频操作放在同一屏”、“焦点规则：先确认“操作”，再确认“语音””、“工单窗口：先记录，再执行呼叫动作”、“线路状态：用可视化面板降低误操作”、“接续流程：从来话到三方沟通”、“转接与交叉操作：语音不一定跟着窗口走”、“合并与拆分：处理复杂通话关系”、“查号：用更少输入找到正确联系人”；第21至第23页依次对应“即时会议：临时协同但不污染会议列表”、“座席互联：用于班组内快速沟通”、“主菜单：按业务动作重新归类”。第24、25页依次对应“运行流程与管理要点”和“谢谢聆听”。",
                lambda p, s: check_plus1_pages_and_order(p, s, file_path),
            ),
            (
                5,
                "PDF转换页面文字：第1至第8页及第24、25页标题、正文、编号和说明均为可编辑文本框。流程框、箭头、圆形、图标、分隔线和示意结构均为可编辑形状或可编辑组合对象。最终第9至第23页的文字、表格、形状、流程对象和图标保持可编辑。",
                check_plus2_editable,
            ),
            (
                5,
                "全部25页背景统一采用浅蓝色或浅蓝渐变背景，不出现与整体风格割裂的深黑、深红或高饱和背景。标题统一使用深蓝色或蓝绿色文字，字体、字号、加粗和左对齐方式保持一致。正文中中文字体、英文字体、字号、行距和项目符号样式统一，正文清晰可读。使用白色、淡蓝色或蓝绿色圆角框，边框颜色、线宽和圆角程度保持一致。",
                check_plus3_style,
            ),
            (
                5,
                "全部25页页码：按照最终合并顺序连续显示1至25。",
                check_plus4_page_numbers,
            ),
        ]

        max_score = sum(pts for pts, _, _ in plus_items)
        result["max_score"] = max_score

        score = 0
        dim2_items = []

        # 加分项
        for pts, rule_text, fn in plus_items:
            ok, _ = fn(prs, slide_data)
            delta = pts if ok else 0
            score += delta
            dim2_items.append({
                "rule": rule_text,
                "max_delta": pts,
                "delta": delta,
                "hit": bool(ok),
                "detail": "",
            })

        result["dim2_items"] = dim2_items
        result["total_score"] = score
        return result

    except Exception as e:
        result["status"] = "error"
        result["error"] = f"{type(e).__name__}: {e}"
        return result


if __name__ == "__main__":
    # 本地调试用：默认使用脚本所在目录
    default_dir = os.path.dirname(os.path.abspath(__file__))
    target_dir = sys.argv[1] if len(sys.argv) > 1 else default_dir
    print(json.dumps(evaluate(target_dir), ensure_ascii=False, indent=2))
