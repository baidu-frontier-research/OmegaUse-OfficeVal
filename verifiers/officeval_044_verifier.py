# -*- coding: utf-8 -*-
"""自动评估脚本：新航路开辟说课PPT_目录视频修改版.pptx"""

import sys
import os
import json
import re
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

SCRIPT_ID = "044"
FILE_NAME = '新航路开辟说课PPT_目录视频修改版.pptx'

ALL_RULES = [
    {"rule": "+1：第8页不再出现旧课程文字", "max_delta": 1},
    {"rule": "+1：第8页'目录'标题格式正确（宋体96号白色、左上半部居中）", "max_delta": 1},
    {"rule": "+1：第8页三个目录条目完整、仿宋44号一致、标题下方均匀居中", "max_delta": 1},
    {"rule": "+1：第8页三个目录条目均设置点击跳转，且分别跳转到第9/16/21页", "max_delta": 1},
    {"rule": "+5：第9页视频为媒体对象、尺寸高12.06宽14.22cm、左侧图片区域、点击触发播放", "max_delta": 5},
    {"rule": "-3：第8页没有被改成目录页", "max_delta": -3},
    {"rule": "-1：第8页目录文字是截图或图片", "max_delta": -1},
    {"rule": "-1：第8页'目录'标题不是宋体或字号不是96号，或位置出现在页面左侧或右侧", "max_delta": -1},
    {"rule": "-1：第8页目录条目字体不是仿宋或字号明显不是44号", "max_delta": -1},
    {"rule": "-1：第8页缺少'第一篇 探索！''第二篇 讨论''第三篇 作业'中的任意一项及以上", "max_delta": -1},
    {"rule": "-3：第8页三个目录条目均未设置点击跳转", "max_delta": -3},
    {"rule": "-1：第9页没有插入视频，或只插入了视频截图", "max_delta": -1},
    {"rule": "-1：第9页视频无法播放，或交付后视频资源丢失", "max_delta": -1},
    {"rule": "-1：第9页视频出现遮挡右侧浅米色文字卡片、顶部导航栏或页面标题'情境导入：一封来自港口的委托信'", "max_delta": -1},
    {"rule": "-1：点击后只显示视频但不能自动播放，需要用户再次手动点击播放键", "max_delta": -1},
    {"rule": "-1：第9页中'情境导入：一封来自港口的委托信'被删除、移动到第8页或第9页以外的主要页面内容", "max_delta": -1},
    {"rule": "-1：第十一页出现明显版式错乱，例如标题'动因一：商业利益与黄金想象'跑出页面、图片遮挡文字、卡片重叠", "max_delta": -1},
    {"rule": "-1：PPT页面尺寸未保持16:9宽屏比例，变成4:3或其他比例", "max_delta": -1},
    {"rule": "-1：'目录'标题与深蓝背景颜色相近，也与图片或目录条目重叠", "max_delta": -1},
    {"rule": "-1：第8页目录标题和三个目录条目均为可编辑文本框，是截图、图片或合成背景", "max_delta": -1},
    {"rule": "-1：目录跳转在放映模式下不可用，点击后跳到空白页、外部网页或错误页面", "max_delta": -1},
    {"rule": "-1：第9页顶部深蓝色导航栏被删除或移动到页面下方，导航栏横向铺满页面顶部", "max_delta": -1},
    {"rule": "-1：第9页顶部'教学理念与立意''教学前端分析''教学过程设计''教学评价'等导航栏文字被删除，文字变为蓝色或接近蓝色，或变成纵向排列", "max_delta": -1},
    {"rule": "-1：第9页标题'情境导入：一封来自港口的委托信'，字体被改为宋体或接近宋体，字号超过26-30号，加粗，颜色改为白色或接近白色的浅色字体", "max_delta": -1},
    {"rule": "-1：第9页右侧浅米色文字卡片被删除，卡片大致位于页面右半部分，高度超过11-13厘米，宽度超过13-15厘米，带浅棕色细边框", "max_delta": -1},
    {"rule": "-1：第9页右侧文字卡片内'致年轻的航海学徒：''港口的货栈越来越拥挤''导入意图：以港口场景激活学生经验，提出本课核心问题'等文字内容被删除", "max_delta": -1},
    {"rule": "-1：第9页右下角页码'09'，被删除或被视频遮挡", "max_delta": -1},
]


def emu_to_cm(emu):
    return emu / 360000


def color_is_white_or_near(rgb_str):
    """判断颜色是否为白色或接近白色"""
    if not rgb_str:
        return False
    r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
    return r > 200 and g > 200 and b > 200


def color_is_blue_or_near(rgb_str):
    """判断颜色是否为蓝色或接近蓝色"""
    if not rgb_str:
        return False
    r, g, b = int(rgb_str[0:2], 16), int(rgb_str[2:4], 16), int(rgb_str[4:6], 16)
    return b > 150 and b > r + 50 and b > g + 50


def get_font_info(para):
    """获取段落第一个run的字体信息"""
    if para.runs:
        r = para.runs[0]
        color_rgb = None
        try:
            if r.font.color and r.font.color.type:
                color_rgb = str(r.font.color.rgb)
        except:
            pass
        return {
            'name': r.font.name,
            'size': r.font.size,
            'bold': r.font.bold,
            'color': color_rgb
        }
    return None


def check_dimension1(prs, file_path):  # file_path 保留以兼容调用方,当前维度一不需要读包内文件
    """维度一：可用与可修改性检查"""
    del file_path
    issues = []

    # 可编辑性检查: 交付物应为"可编辑PPT",而不是把每页拍平成一张图片的"截图版"。
    # 判定思路: 逐页扫描,若某页只由整页尺寸的图片构成(近乎覆盖全页且没有任何
    # 可编辑文本run),视为整页图片化,该页不再可编辑。
    # 仅当"存在整页图片化"的页面时判为不通过;正常演示页面即便含大背景图,
    # 只要其上还留有可编辑文字(标题/正文/页码)就算可编辑。
    slide_w_emu = prs.slide_width or 0
    slide_h_emu = prs.slide_height or 0
    for i, sl in enumerate(prs.slides):
        has_editable_text = any(
            s.has_text_frame and any(p.text.strip() for p in s.text_frame.paragraphs)
            for s in sl.shapes
        )
        if has_editable_text:
            continue
        # 页面上没有任何可编辑文本run: 若同时存在近乎覆盖全页的图片,视为截图化
        has_fullpage_pic = any(
            s.shape_type == MSO_SHAPE_TYPE.PICTURE
            and s.width and s.height
            and slide_w_emu and slide_h_emu
            and s.width >= slide_w_emu * 0.9
            and s.height >= slide_h_emu * 0.9
            for s in sl.shapes
        )
        if has_fullpage_pic:
            issues.append(f"第{i+1}页整页图片化,文件不可编辑")
            break

    blank_count = 0
    for i, slide in enumerate(prs.slides):
        has_content = False
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                has_content = True
                break
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_content = True
                break
            if 'MEDIA' in str(shape.shape_type):
                has_content = True
                break
        if not has_content:
            blank_count += 1
            if blank_count >= 2:
                issues.append(f"第{i}页和第{i+1}页为连续空白页")
                break
        else:
            blank_count = 0

    return issues


def check_dimension2(prs):
    """维度二：完成度评分"""
    hits = []
    slide8 = prs.slides[7]
    slide9 = prs.slides[8]

    slide8_texts = []
    slide8_shapes_with_text = []
    for shape in slide8.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                slide8_texts.append(text)
                slide8_shapes_with_text.append(shape)

    all_slide8_text = ' '.join(slide8_texts)

    old_texts = ["第6课 第一子目", "新航路开辟的动因", "历史课堂说课设计", "澜川市星商中学 · 林沐辰"]

    def _normalize_text(s):
        s = re.sub(r'\s+', '', s)
        for dot in ('·', '•', '‧', '・', '⋅', '∙', '·'):
            s = s.replace(dot, '')
        return s

    norm_slide8 = _normalize_text(all_slide8_text)
    has_old_text = any(_normalize_text(ot) in norm_slide8 for ot in old_texts)
    if not has_old_text:
        hits.append((1, "+1：第8页不再出现旧课程文字"))

    mulu_ok = False
    for shape in slide8_shapes_with_text:
        text = shape.text_frame.text.strip()
        if text != "目录":
            continue
        para = shape.text_frame.paragraphs[0]
        runs = para.runs
        if not runs:
            continue
        is_song = all(r.font.name and '宋' in r.font.name for r in runs)
        is_96 = all(r.font.size and abs(r.font.size - Pt(96)) < Pt(2) for r in runs)
        def _run_color(r):
            try:
                if r.font.color and r.font.color.type:
                    return str(r.font.color.rgb)
            except Exception:
                pass
            return None
        is_white = all(color_is_white_or_near(_run_color(r)) for r in runs)
        shape_v_center = shape.top + shape.height // 2
        in_top_half = shape_v_center < prs.slide_height // 2
        at_left = shape.left < prs.slide_width // 2
        is_centered = para.alignment == PP_ALIGN.CENTER
        if is_song and is_96 and is_white and in_top_half and at_left and is_centered:
            mulu_ok = True
    if mulu_ok:
        hits.append((1, "+1：第8页'目录'标题格式正确（宋体96号白色、左上半部居中）"))

    entries_required = ["第一篇 探索！", "第二篇 讨论", "第三篇 作业"]

    def _norm_space(s):
        return re.sub(r'\s+', '', s)

    entry_shapes = {}
    for shape in slide8_shapes_with_text:
        text_norm = _norm_space(shape.text_frame.text.strip())
        for entry in entries_required:
            if _norm_space(entry) == text_norm:
                entry_shapes[entry] = shape

    all_entries_present = all(e in entry_shapes for e in entries_required)

    entries_font_ok = True
    seen_fonts = set()
    seen_sizes = set()
    for entry in entries_required:
        shape = entry_shapes.get(entry)
        if not shape:
            entries_font_ok = False
            continue
        para = shape.text_frame.paragraphs[0]
        runs = para.runs
        if not runs:
            entries_font_ok = False
            continue
        if not all(r.font.name and '仿宋' in r.font.name for r in runs):
            entries_font_ok = False
        if not all(r.font.size and abs(r.font.size - Pt(44)) < Pt(2) for r in runs):
            entries_font_ok = False
        for r in runs:
            if r.font.name:
                seen_fonts.add(r.font.name)
            if r.font.size:
                seen_sizes.add(r.font.size)
    consistent = len(seen_fonts) == 1 and len(seen_sizes) == 1

    title_bottom = None
    for shape in slide8_shapes_with_text:
        if shape.text_frame.text.strip() == "目录":
            title_bottom = shape.top + shape.height
            break
    below_title = True
    if title_bottom is not None and all_entries_present:
        for entry in entries_required:
            if entry_shapes[entry].top < title_bottom:
                below_title = False
    else:
        below_title = all_entries_present

    spacing_even = True
    if all_entries_present:
        centers = sorted(entry_shapes[e].top + entry_shapes[e].height // 2
                         for e in entries_required)
        gaps = [centers[1] - centers[0], centers[2] - centers[1]]
        avg_gap = sum(gaps) / len(gaps)
        if avg_gap <= 0 or abs(gaps[0] - gaps[1]) > avg_gap * 0.2:
            spacing_even = False

    centered_ok = True
    if all_entries_present:
        page_cx = prs.slide_width // 2
        for entry in entries_required:
            shape = entry_shapes[entry]
            para = shape.text_frame.paragraphs[0]
            if para.alignment != PP_ALIGN.CENTER:
                centered_ok = False
            shape_cx = shape.left + shape.width // 2
            if abs(shape_cx - page_cx) > prs.slide_width * 0.1:
                centered_ok = False
    else:
        centered_ok = False

    entries_found = [e for e in entries_required if e in entry_shapes]

    if all_entries_present and entries_font_ok and consistent and below_title and spacing_even and centered_ok:
        hits.append((1, "+1：第8页三个目录条目完整、仿宋44号一致、标题下方均匀居中"))

    slide_display_page = {}
    for idx, sl in enumerate(prs.slides):
        slide_display_page[sl.part.partname] = idx + 1

    expected_jumps = {"第一篇 探索！": 9, "第二篇 讨论": 16, "第三篇 作业": 21}

    def _norm_space2(s):
        return re.sub(r'\s+', '', s)

    jump_info = {}
    for shape in slide8_shapes_with_text:
        text_norm = _norm_space2(shape.text_frame.text.strip())
        for entry in expected_jumps:
            if _norm_space2(entry) != text_norm:
                continue
            xml_str = shape._element.xml
            has_action = ('hlinkClick' in xml_str and 'hlinksldjump' in xml_str)
            page = None
            if has_action:
                rids = re.findall(r'r:id="(rId\d+)"', xml_str)
                for rid in rids:
                    if rid in slide8.part.rels:
                        rel = slide8.part.rels[rid]
                        if not rel.is_external:
                            try:
                                tgt_part = rel.target_part
                                if tgt_part.partname in slide_display_page:
                                    page = slide_display_page[tgt_part.partname]
                            except Exception:
                                page = None
            jump_info[entry] = {'has_action': has_action, 'page': page}

    all_jumps_correct = True
    for entry, want_page in expected_jumps.items():
        info = jump_info.get(entry)
        if not info or not info['has_action'] or info['page'] != want_page:
            all_jumps_correct = False

    if all_jumps_correct:
        hits.append((1, "+1：第8页三个目录条目均设置点击跳转，且分别跳转到第9/16/21页"))

    video_score_5 = True
    video_shape = None
    for shape in slide9.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.MEDIA:
            video_shape = shape
            break

    if video_shape is None:
        video_score_5 = False
    else:
        # rubric 尺寸: 宽12.06 cm、高14.22 cm(误差 ±0.3 cm)
        w_cm = emu_to_cm(video_shape.width)
        h_cm = emu_to_cm(video_shape.height)
        width_ok = abs(w_cm - 12.06) <= 0.3
        height_ok = abs(h_cm - 14.22) <= 0.3
        if not (width_ok and height_ok):
            video_score_5 = False
        # 位于第9页"左侧图片区域": 视频矩形须实质位于页面左半区图片区域内。
        # 判据:
        #   (a) 视频矩形整体落在页面左半页(允许小容差),不能越过垂直中线
        #   (b) 与"左半页矩形"的重叠面积占视频自身面积的比例达到大部分(≥0.9),
        #       等价于办公软件里肉眼可见"视频完全在左半图片区"
        slide_w = prs.slide_width or 0
        slide_h = prs.slide_height or 0
        vx, vy = video_shape.left or 0, video_shape.top or 0
        vw, vh = video_shape.width or 0, video_shape.height or 0
        tol = Cm(0.5)
        left_area = (0, 0, slide_w // 2, slide_h)
        ox = max(0, min(vx + vw, left_area[2]) - max(vx, left_area[0]))
        oy = max(0, min(vy + vh, left_area[3]) - max(vy, left_area[1]))
        overlap = ox * oy
        video_area = vw * vh
        in_left_area = (
            vx >= -tol and vy >= -tol
            and vx + vw <= slide_w // 2 + tol
            and vy + vh <= slide_h + tol
            and video_area > 0
            and overlap >= video_area * 0.9
        )
        if not in_left_area:
            video_score_5 = False
        slide_xml = slide9._element.xml
        timing_match = re.search(r'<p:timing.*?</p:timing>', slide_xml, re.DOTALL)
        if timing_match:
            timing_xml = timing_match.group(0)
            if 'delay="indefinite"' not in timing_xml:
                video_score_5 = False
            if 'afterPrevious' in timing_xml or 'withPrevious' in timing_xml:
                video_score_5 = False
        else:
            video_score_5 = False

    if video_score_5:
        hits.append((5, "+5：第9页视频为媒体对象、尺寸宽12.06高14.22cm、左侧图片区域、点击触发播放"))

    # === 扣分项 ===
    norm_slide8_for_mulu = re.sub(r'\s+', '', all_slide8_text)
    if "目录" not in norm_slide8_for_mulu:
        hits.append((-3, "-3：第8页没有被改成目录页"))

    mulu_in_editable_text = False
    for shape in slide8_shapes_with_text:
        if "目录" in re.sub(r'\s+', '', shape.text_frame.text):
            mulu_in_editable_text = True
            break
    if not mulu_in_editable_text:
        hits.append((-1, "-1：第8页目录文字是截图或图片"))

    for shape in slide8_shapes_with_text:
        text = shape.text_frame.text.strip()
        if text == "目录":
            runs = shape.text_frame.paragraphs[0].runs
            is_song = bool(runs) and all(r.font.name and '宋' in r.font.name for r in runs)
            is_96 = bool(runs) and all(r.font.size and abs(r.font.size - Pt(96)) < Pt(2) for r in runs)
            shape_cx = shape.left + shape.width // 2
            page_cx = prs.slide_width // 2
            off_center = abs(shape_cx - page_cx) > prs.slide_width * 0.1
            if (not is_song) or (not is_96) or off_center:
                hits.append((-1, "-1：第8页'目录'标题不是宋体或字号不是96号，或位置出现在页面左侧或右侧"))
            break

    entries_font_bad = False
    for entry in entries_required:
        shape = entry_shapes.get(entry)
        if not shape:
            continue
        runs = shape.text_frame.paragraphs[0].runs
        if not runs:
            entries_font_bad = True
            continue
        if not all(r.font.name and '仿宋' in r.font.name for r in runs):
            entries_font_bad = True
        if not all(r.font.size and abs(r.font.size - Pt(44)) <= Pt(4) for r in runs):
            entries_font_bad = True
    if entries_font_bad:
        hits.append((-1, "-1：第8页目录条目字体不是仿宋或字号明显不是44号"))

    slide8_all_norm = re.sub(r'\s+', '', all_slide8_text)
    missing_any_entry = any(
        re.sub(r'\s+', '', entry) not in slide8_all_norm
        for entry in entries_required
    )
    if missing_any_entry:
        hits.append((-1, "-1：第8页缺少'第一篇 探索！''第二篇 讨论''第三篇 作业'中的任意一项及以上"))

    entries_with_jump = 0
    for shape in slide8_shapes_with_text:
        text_norm = re.sub(r'\s+', '', shape.text_frame.text.strip())
        if text_norm not in {re.sub(r'\s+', '', e) for e in entries_required}:
            continue
        xml_str = shape._element.xml
        if 'hlinkClick' in xml_str and 'hlinksldjump' in xml_str:
            entries_with_jump += 1
    if entries_with_jump == 0:
        hits.append((-3, "-3：第8页三个目录条目均未设置点击跳转"))

    has_video_media = any(s.shape_type == MSO_SHAPE_TYPE.MEDIA for s in slide9.shapes)
    if not has_video_media:
        hits.append((-1, "-1：第9页没有插入视频，或只插入了视频截图"))

    video_media_shape = None
    for s in slide9.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.MEDIA:
            video_media_shape = s
            break

    if video_media_shape is not None:
        video_ok = False
        xml_str = video_media_shape._element.xml
        rids = set(re.findall(r'r:(?:link|embed|id)="(rId\d+)"', xml_str))
        for rid in rids:
            if rid in slide9.part.rels:
                rel = slide9.part.rels[rid]
                if rel.reltype.endswith('/video') and not rel.is_external:
                    try:
                        blob = rel.target_part.blob
                        if blob and len(blob) > 0:
                            video_ok = True
                    except Exception:
                        video_ok = False
        if not video_ok:
            hits.append((-1, "-1：第9页视频无法播放，或交付后视频资源丢失"))

    vshape = None
    for s in slide9.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.MEDIA:
            vshape = s
            break
    if vshape is not None:
        tol = Cm(0.1)
        def _rects_overlap(a_left, a_top, a_w, a_h, b_left, b_top, b_w, b_h):
            ix = min(a_left + a_w, b_left + b_w) - max(a_left, b_left)
            iy = min(a_top + a_h, b_top + b_h) - max(a_top, b_top)
            return ix > tol and iy > tol
        vx, vy, vw, vh = vshape.left, vshape.top, vshape.width, vshape.height
        occlude = False
        for s in slide9.shapes:
            if s is vshape:
                continue
            is_title = (s.has_text_frame and "情境导入" in s.text_frame.text)
            is_navbar = (s.top < Cm(0.1) and
                         s.width > prs.slide_width * 0.8 and
                         s.height < Cm(2) and s.height > 0)
            is_right_card = ((s.left + s.width // 2) > prs.slide_width // 2 and
                             s.width > Cm(10) and s.height > Cm(8))
            if not (is_title or is_navbar or is_right_card):
                continue
            if _rects_overlap(vx, vy, vw, vh, s.left, s.top, s.width, s.height):
                occlude = True
                break
        if occlude:
            hits.append((-1, "-1：第9页视频出现遮挡右侧浅米色文字卡片、顶部导航栏或页面标题'情境导入：一封来自港口的委托信'"))

    vshape = None
    for s in slide9.shapes:
        if s.shape_type == MSO_SHAPE_TYPE.MEDIA:
            vshape = s
            break
    if vshape is not None:
        slide_xml = slide9._element.xml
        timing_match = re.search(r'<p:timing.*?</p:timing>', slide_xml, re.DOTALL)
        has_video_play_action = False
        if timing_match:
            timing_xml = timing_match.group(0)
            target_spid = str(vshape.shape_id)
            for vnode in re.finditer(r'<p:video>.*?</p:video>', timing_xml, re.DOTALL):
                block = vnode.group(0)
                sp = re.search(r'<p:spTgt\s+spid="(\d+)"', block)
                if sp and sp.group(1) == target_spid:
                    has_video_play_action = True
                    break
        if not has_video_play_action:
            hits.append((-1, "-1：点击后只显示视频但不能自动播放，需要用户再次手动点击播放键"))

    title_full = "情境导入：一封来自港口的委托信"
    title_full_norm = re.sub(r'\s+', '', title_full)

    def _para_norm_texts(shape):
        if not shape.has_text_frame:
            return []
        out = []
        for para in shape.text_frame.paragraphs:
            t = re.sub(r'\s+', '', ''.join(r.text for r in para.runs))
            if t:
                out.append(t)
        return out

    title_on_slide9 = False
    for shape in slide9.shapes:
        if any(title_full_norm in t for t in _para_norm_texts(shape)):
            title_on_slide9 = True
            break

    title_elsewhere = False
    for idx, sl in enumerate(prs.slides):
        if idx == 8:
            continue
        for shape in sl.shapes:
            for t in _para_norm_texts(shape):
                if t == title_full_norm:
                    title_elsewhere = True
                    break
            if title_elsewhere:
                break
        if title_elsewhere:
            break

    if (not title_on_slide9) or title_elsewhere:
        hits.append((-1, "-1：第9页中'情境导入：一封来自港口的委托信'被删除、移动到第8页或第9页以外的主要页面内容"))

    slide11 = prs.slides[10]
    tol_out = Cm(0.3)
    layout_bad = False

    def _rect(shape):
        return (shape.left, shape.top, shape.width, shape.height)

    def _overlap_frac(a, b):
        ax, ay, aw, ah = a
        bx, by, bw, bh = b
        if aw <= 0 or ah <= 0 or bw <= 0 or bh <= 0:
            return 0.0
        ix = min(ax + aw, bx + bw) - max(ax, bx)
        iy = min(ay + ah, by + bh) - max(ay, by)
        if ix <= 0 or iy <= 0:
            return 0.0
        inter = ix * iy
        smaller = min(aw * ah, bw * bh)
        return inter / smaller if smaller > 0 else 0.0

    for shape in slide11.shapes:
        if not (shape.has_text_frame and shape.text_frame.text.strip()):
            continue
        if (shape.left < -tol_out or shape.top < -tol_out or
            shape.left + shape.width > prs.slide_width + tol_out or
            shape.top + shape.height > prs.slide_height + tol_out):
            layout_bad = True
            break

    if not layout_bad:
        pics = [s for s in slide11.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        texts = [s for s in slide11.shapes
                 if s.has_text_frame and s.text_frame.text.strip()]
        for pic in pics:
            for txt in texts:
                if _overlap_frac(_rect(pic), _rect(txt)) > 0.15:
                    layout_bad = True
                    break
            if layout_bad:
                break

    if not layout_bad:
        cards = [s for s in slide11.shapes
                 if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
                 and s.width > Cm(4) and s.height > Cm(3)]
        for i in range(len(cards)):
            for j in range(i + 1, len(cards)):
                if _overlap_frac(_rect(cards[i]), _rect(cards[j])) > 0.15:
                    layout_bad = True
                    break
            if layout_bad:
                break

    if layout_bad:
        hits.append((-1, "-1：第十一页出现明显版式错乱，例如标题'动因一：商业利益与黄金想象'跑出页面、图片遮挡文字、卡片重叠"))

    ratio = prs.slide_width / prs.slide_height
    if abs(ratio - 16 / 9) / (16 / 9) > 0.01:
        hits.append((-1, "-1：PPT页面尺寸未保持16:9宽屏比例，变成4:3或其他比例"))

    mulu_title_shape = None
    for shape in slide8_shapes_with_text:
        if shape.text_frame.text.strip() == "目录":
            mulu_title_shape = shape
            break

    if mulu_title_shape is not None:
        title_rgb = None
        for r in mulu_title_shape.text_frame.paragraphs[0].runs:
            try:
                if r.font.color and r.font.color.type:
                    title_rgb = str(r.font.color.rgb)
                    break
            except Exception:
                pass

        bg_rgb = None
        page_area = prs.slide_width * prs.slide_height
        for s in slide8.shapes:
            if s.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
            if s.width * s.height < page_area * 0.9:
                continue
            try:
                if s.fill.type == 1:
                    rgb_str = str(s.fill.fore_color.rgb)
                    if color_is_blue_or_near(rgb_str):
                        bg_rgb = rgb_str
                        break
            except Exception:
                continue

        color_similar = False
        if title_rgb and bg_rgb:
            tr, tg, tb = int(title_rgb[0:2], 16), int(title_rgb[2:4], 16), int(title_rgb[4:6], 16)
            br, bg_, bb = int(bg_rgb[0:2], 16), int(bg_rgb[2:4], 16), int(bg_rgb[4:6], 16)
            dist = ((tr - br) ** 2 + (tg - bg_) ** 2 + (tb - bb) ** 2) ** 0.5
            if dist <= 80:
                color_similar = True

        def _rect2(s):
            return (s.left, s.top, s.width, s.height)

        def _overlap_frac2(a, b):
            ax, ay, aw, ah = a
            bx, by, bw, bh = b
            if aw <= 0 or ah <= 0 or bw <= 0 or bh <= 0:
                return 0.0
            ix = min(ax + aw, bx + bw) - max(ax, bx)
            iy = min(ay + ah, by + bh) - max(ay, by)
            if ix <= 0 or iy <= 0:
                return 0.0
            smaller = min(aw * ah, bw * bh)
            return (ix * iy) / smaller if smaller > 0 else 0.0

        entry_norms = {re.sub(r'\s+', '', e) for e in entries_required}
        overlap_hit = False
        title_rect = _rect2(mulu_title_shape)
        for s in slide8.shapes:
            if s is mulu_title_shape:
                continue
            is_pic = (s.shape_type == MSO_SHAPE_TYPE.PICTURE)
            is_entry = (s.has_text_frame and
                        re.sub(r'\s+', '', s.text_frame.text.strip()) in entry_norms)
            if not (is_pic or is_entry):
                continue
            if _overlap_frac2(title_rect, _rect2(s)) > 0.15:
                overlap_hit = True
                break

        if color_similar and overlap_hit:
            hits.append((-1, "-1：'目录'标题与深蓝背景颜色相近，也与图片或目录条目重叠"))

    def _text_in_editable_box(target_norm):
        for s in slide8.shapes:
            if not s.has_text_frame:
                continue
            for para in s.text_frame.paragraphs:
                if target_norm in re.sub(r'\s+', '', ''.join(r.text for r in para.runs)):
                    return True
        return False

    title_norm = re.sub(r'\s+', '', "目录")
    targets_norm = [title_norm] + [re.sub(r'\s+', '', e) for e in entries_required]
    any_not_editable = any(not _text_in_editable_box(t) for t in targets_norm)
    if any_not_editable:
        hits.append((-1, "-1：第8页目录标题和三个目录条目均为可编辑文本框，是截图、图片或合成背景"))

    display_page_map2 = {}
    for i2, sl2 in enumerate(prs.slides):
        display_page_map2[sl2.part.partname] = i2 + 1
    expected_pages2 = {re.sub(r'\s+', '', k): v
                       for k, v in {"第一篇 探索！": 9,
                                     "第二篇 讨论": 16,
                                     "第三篇 作业": 21}.items()}

    def _slide_is_blank(sl):
        for s in sl.shapes:
            if s.has_text_frame and s.text_frame.text.strip():
                return False
            if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return False
            if s.shape_type == MSO_SHAPE_TYPE.MEDIA:
                return False
        return True

    partname_to_slide = {sl.part.partname: sl for sl in prs.slides}

    jump_broken = False
    for shape in slide8_shapes_with_text:
        text_norm = re.sub(r'\s+', '', shape.text_frame.text.strip())
        if text_norm not in expected_pages2:
            continue
        xml_str = shape._element.xml
        if 'hlinkClick' not in xml_str:
            continue
        m = re.search(r'<a:hlinkClick[^>]*?r:id="(rId\d+)"', xml_str)
        if not m:
            jump_broken = True
            continue
        rid = m.group(1)
        if rid not in slide8.part.rels:
            jump_broken = True
            continue
        rel = slide8.part.rels[rid]
        if rel.is_external:
            jump_broken = True
            continue
        try:
            tgt_part = rel.target_part
        except Exception:
            jump_broken = True
            continue
        if tgt_part.partname not in partname_to_slide:
            jump_broken = True
            continue
        target_page = display_page_map2.get(tgt_part.partname)
        target_slide = partname_to_slide[tgt_part.partname]
        if _slide_is_blank(target_slide):
            jump_broken = True
            continue
        if target_page != expected_pages2[text_norm]:
            jump_broken = True
            continue

    if jump_broken:
        hits.append((-1, "-1：目录跳转在放映模式下不可用，点击后跳到空白页、外部网页或错误页面"))

    nav_top_ok = False
    nav_moved_down = False
    for shape in slide9.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if shape.width < prs.slide_width * 0.8:
            continue
        if shape.height <= 0 or shape.height > Cm(3):
            continue
        try:
            if shape.fill.type != 1:
                continue
            rgb_str = str(shape.fill.fore_color.rgb)
        except Exception:
            continue
        try:
            r_v = int(rgb_str[0:2], 16)
            g_v = int(rgb_str[2:4], 16)
            b_v = int(rgb_str[4:6], 16)
        except Exception:
            continue
        is_dark_blue = (max(r_v, g_v, b_v) < 150 and
                        b_v > r_v + 10 and b_v > g_v + 10)
        if not is_dark_blue:
            continue
        if shape.top < prs.slide_height * 0.1:
            nav_top_ok = True
        elif shape.top > prs.slide_height * 0.5:
            nav_moved_down = True

    if (not nav_top_ok) or nav_moved_down:
        hits.append((-1, "-1：第9页顶部深蓝色导航栏被删除或移动到页面下方，导航栏横向铺满页面顶部"))

    nav_texts = ["教学理念与立意", "教学前端分析", "教学过程设计", "教学评价"]
    nav_norms = {re.sub(r'\s+', '', t): t for t in nav_texts}

    nav_top_boundary = prs.slide_height * 0.1

    found_nav_shapes = {}
    for shape in slide9.shapes:
        if not shape.has_text_frame:
            continue
        if shape.top >= nav_top_boundary:
            continue
        text_norm = re.sub(r'\s+', '', shape.text_frame.text.strip())
        if text_norm in nav_norms and text_norm not in found_nav_shapes:
            found_nav_shapes[text_norm] = shape

    missing_nav = len(found_nav_shapes) < len(nav_texts)

    color_blue = False
    vertical = False
    for norm_text, shape in found_nav_shapes.items():
        for para in shape.text_frame.paragraphs:
            for r in para.runs:
                try:
                    if r.font.color and r.font.color.type:
                        rgb_str = str(r.font.color.rgb)
                        if color_is_blue_or_near(rgb_str):
                            color_blue = True
                except Exception:
                    pass
        xml_s = shape._element.xml
        m = re.search(r'<a:bodyPr[^>]*vert="([^"]+)"', xml_s)
        if m and m.group(1) in ('eaVert', 'vert', 'vert270', 'mongolianVert'):
            vertical = True
        else:
            if shape.height > shape.width * 2:
                para_count = sum(
                    1 for para in shape.text_frame.paragraphs
                    if ''.join(rr.text for rr in para.runs).strip()
                )
                if para_count >= max(2, len(norm_text) - 1):
                    vertical = True

    if missing_nav or color_blue or vertical:
        hits.append((-1, "-1：第9页顶部'教学理念与立意''教学前端分析''教学过程设计''教学评价'等导航栏文字被删除，文字变为蓝色或接近蓝色，或变成纵向排列"))

    title_full2 = "情境导入：一封来自港口的委托信"
    title_full2_norm = re.sub(r'\s+', '', title_full2)

    title_shape = None
    for shape in slide9.shapes:
        if not shape.has_text_frame:
            continue
        norm = re.sub(r'\s+', '', shape.text_frame.text.strip())
        if norm == title_full2_norm:
            title_shape = shape
            break

    if title_shape is not None:
        is_song_like = False
        oversize = False
        is_bold = False
        is_white_like = False
        for para in title_shape.text_frame.paragraphs:
            for r in para.runs:
                if r.font.name and '宋' in r.font.name:
                    is_song_like = True
                if r.font.size and r.font.size > Pt(30):
                    oversize = True
                if r.font.bold:
                    is_bold = True
                try:
                    if r.font.color and r.font.color.type:
                        if color_is_white_or_near(str(r.font.color.rgb)):
                            is_white_like = True
                except Exception:
                    pass
        if is_song_like and oversize and is_bold and is_white_like:
            hits.append((-1, "-1：第9页标题'情境导入：一封来自港口的委托信'，字体被改为宋体或接近宋体，字号超过26-30号，加粗，颜色改为白色或接近白色的浅色字体"))

    def _rgb_hex_to_tuple(s):
        try:
            return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
        except Exception:
            return None

    def _is_light_cream(rgb_str):
        t = _rgb_hex_to_tuple(rgb_str)
        if not t:
            return False
        r, g, b = t
        return (r >= g >= b and min(r, g, b) >= 200
                and not (r >= 245 and g >= 245 and b >= 245))

    def _is_light_brown(rgb_str):
        t = _rgb_hex_to_tuple(rgb_str)
        if not t:
            return False
        r, g, b = t
        return r > g > b and r >= 150 and b >= 100

    card_found = False
    for shape in slide9.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            continue
        if (shape.left + shape.width // 2) <= prs.slide_width // 2:
            continue
        if emu_to_cm(shape.height) <= 11 or emu_to_cm(shape.width) <= 13:
            continue
        try:
            if shape.fill.type != 1:
                continue
            fill_rgb = str(shape.fill.fore_color.rgb)
        except Exception:
            continue
        if not _is_light_cream(fill_rgb):
            continue
        xml_s = shape._element.xml
        has_light_brown_thin_border = False
        for m in re.finditer(r'<a:ln([^>/]*)>(.*?)</a:ln>', xml_s, re.DOTALL):
            attrs, body = m.group(1), m.group(2)
            if '<a:noFill' in body:
                continue
            w_m = re.search(r'w="(\d+)"', attrs)
            width_emu = int(w_m.group(1)) if w_m else 0
            if width_emu > 38100:
                continue
            color_m = re.search(r'<a:srgbClr\s+val="([0-9A-Fa-f]{6})"', body)
            if not color_m:
                continue
            if _is_light_brown(color_m.group(1)):
                has_light_brown_thin_border = True
                break
        if not has_light_brown_thin_border:
            continue
        card_found = True
        break

    if not card_found:
        hits.append((-1, "-1：第9页右侧浅米色文字卡片被删除，卡片大致位于页面右半部分，高度超过11-13厘米，宽度超过13-15厘米，带浅棕色细边框"))

    card_texts = [
        "致年轻的航海学徒：",
        "港口的货栈越来越拥挤",
        "导入意图：以港口场景激活学生经验，提出本课核心问题",
    ]
    card_norms = [re.sub(r'\s+', '', t) for t in card_texts]

    right_text_blob_norm = ''
    for shape in slide9.shapes:
        if not shape.has_text_frame:
            continue
        if (shape.left + shape.width // 2) <= prs.slide_width // 2:
            continue
        right_text_blob_norm += re.sub(r'\s+', '', shape.text_frame.text)

    if any(cn not in right_text_blob_norm for cn in card_norms):
        hits.append((-1, "-1：第9页右侧文字卡片内'致年轻的航海学徒：''港口的货栈越来越拥挤''导入意图：以港口场景激活学生经验，提出本课核心问题'等文字内容被删除"))

    page_num_shape = None
    for shape in slide9.shapes:
        if not shape.has_text_frame:
            continue
        if re.sub(r'\s+', '', shape.text_frame.text) != "09":
            continue
        cx = shape.left + shape.width // 2
        cy = shape.top + shape.height // 2
        if cx > prs.slide_width // 2 and cy > prs.slide_height // 2:
            page_num_shape = shape
            break

    bad_page_num = False
    if page_num_shape is None:
        bad_page_num = True
    else:
        vshape_pn = None
        for s in slide9.shapes:
            if s.shape_type == MSO_SHAPE_TYPE.MEDIA:
                vshape_pn = s
                break
        if vshape_pn is not None:
            ix = min(page_num_shape.left + page_num_shape.width,
                     vshape_pn.left + vshape_pn.width) - max(page_num_shape.left, vshape_pn.left)
            iy = min(page_num_shape.top + page_num_shape.height,
                     vshape_pn.top + vshape_pn.height) - max(page_num_shape.top, vshape_pn.top)
            if ix > 0 and iy > 0:
                pn_area = page_num_shape.width * page_num_shape.height
                if pn_area > 0 and (ix * iy) / pn_area > 0.15:
                    bad_page_num = True

    if bad_page_num:
        hits.append((-1, "-1：第9页右下角页码'09'，被删除或被视频遮挡"))

    return hits


def evaluate(dir_path: str) -> dict:
    """统一评估入口"""
    file_path = os.path.join(dir_path, FILE_NAME)
    result = {
        "id": SCRIPT_ID,
        "file_name": FILE_NAME,
        "status": "ok",
        "error": None,
        "dim1_pass": True,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": 9,
    }

    try:
        if not os.path.isfile(file_path):
            result["status"] = "error"
            result["error"] = f"文件不存在: {FILE_NAME}"
            return result

        prs = Presentation(file_path)

        # 维度一
        dim1_issues = check_dimension1(prs, file_path)
        if dim1_issues:
            result["dim1_pass"] = False
            result["dim1_reason"] = "; ".join(dim1_issues)
            result["dim2_items"] = []
            result["total_score"] = 0
            return result

        # 维度二
        hits = check_dimension2(prs)
        hit_descs = {desc for _, desc in hits}

        dim2_items = []
        for rule_def in ALL_RULES:
            rule_text = str(rule_def["rule"])
            max_delta = int(rule_def["max_delta"])
            hit = rule_text in hit_descs
            delta = max_delta if hit else 0
            # detail 字段按需求统一返回空字符串，不影响命中判定与分数计算
            detail = ""
            dim2_items.append({
                "rule": rule_text,
                "max_delta": max_delta,
                "delta": delta,
                "hit": hit,
                "detail": detail,
            })

        total = sum(item["delta"] for item in dim2_items)
        result["dim2_items"] = dim2_items
        result["total_score"] = total

    except Exception as e:
        result["status"] = "error"
        result["error"] = str(e)

    return result


if __name__ == "__main__":
    dir_path = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    result = evaluate(dir_path)
    print(json.dumps(result, ensure_ascii=False, indent=2))
