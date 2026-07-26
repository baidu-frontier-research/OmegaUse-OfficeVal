# -*- coding: utf-8 -*-
"""
对 "晨岚咖啡品牌运营研究_标题栏与结束页_完成版.pptx" 的自动评估脚本。

评估流程：
  维度1（可用与可修改性）  —— 任意一项不达标，总分直接判为 0，不再检查维度2。
  维度2（完成度）          —— 包含若干加分细则与扣分细则。
                             加分细则：必须满足该细则内的"全部要点"才计这条加分。
                             扣分细则：只要满足该细则内的"任意一个要点"即计这条扣分。

依赖：python-pptx
"""

from __future__ import annotations

import os
import sys
import json
import zipfile
import traceback
from typing import List, Optional, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

# 脚本编号（与文件名 officeval_XXX_verifier.py 保持一致）
SCRIPT_ID = "062"

# ---------- 单位换算 ----------
EMU_PER_CM = 360000
EMU_PER_PT = 12700


def emu_to_cm(v: int) -> float:
    return v / EMU_PER_CM


def emu_to_pt(v: int) -> float:
    return v / EMU_PER_PT


# ---------- 维度2 评分记录器 ----------
class Dim2Log:
    """维度2 评分记录器：每条评分项无论是否命中都会追加一条记录。

    每条记录形如：
        {"rule": "...", "max_delta": 5, "delta": 5|0, "hit": bool, "detail": ""}
    """

    def __init__(self):
        self.items: List[dict] = []

    def record(self, rule: str, max_delta: int, hit: bool, detail: str = ""):
        delta = max_delta if hit else 0
        # 按要求：返回结构中 detail 字段一律置空（不影响命中判定与打分逻辑）
        self.items.append({
            "rule": rule,
            "max_delta": max_delta,
            "delta": delta,
            "hit": bool(hit),
            "detail": "",
        })

    def total(self) -> int:
        return sum(it["delta"] for it in self.items)


# ---------- 颜色相关工具 ----------
def _rgb_of(run_font) -> Optional[Tuple[int, int, int]]:
    """安全地取出 run.font.color.rgb，没有就返回 None。"""
    try:
        c = run_font.color
        if c is None or c.type is None:
            return None
        rgb = c.rgb
        if rgb is None:
            return None
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


def is_blue(rgb: Optional[Tuple[int, int, int]]) -> bool:
    """蓝色判定：B 通道为主，且 B>R 较多。"""
    if rgb is None:
        return False
    r, g, b = rgb
    return b >= 120 and b > r + 20 and b >= g - 20


def is_light_blue_or_white(rgb: Optional[Tuple[int, int, int]]) -> bool:
    if rgb is None:
        return False
    r, g, b = rgb
    if r >= 230 and g >= 230 and b >= 230:
        return True  # 白
    # 浅蓝：B 较高，整体偏亮
    return b >= 180 and (r + g) >= 300 and b >= r


def is_light_blue(rgb: Optional[Tuple[int, int, int]]) -> bool:
    """浅蓝色（不含白色）：B 较高，整体偏亮，且 B > R（去掉白色分支）。"""
    if rgb is None:
        return False
    r, g, b = rgb
    # 排除白色 / 灰色
    if r >= 230 and g >= 230 and b >= 230:
        return False
    return b >= 180 and (r + g) >= 300 and b >= r and b > r + 5


def is_blue_or_light_blue(rgb: Optional[Tuple[int, int, int]]) -> bool:
    """严格"蓝色 或 浅蓝色"：不含白色、不含灰色。"""
    return is_blue(rgb) or is_light_blue(rgb)


# ---------- 非主题色（白/浅/灰）判定 ----------
# 说明：以下若干判定用于"非主题色"色值要求（例如规范中的"白色"、"浅色"、"灰色"等）。
# 对于这类要求，只需符合颜色本身的定义即可通过，不再要求与主题蓝色相近。
def is_white(rgb: Optional[Tuple[int, int, int]]) -> bool:
    """白色（非主题色）：三通道均较高。"""
    if rgb is None:
        return False
    r, g, b = rgb
    return r >= 230 and g >= 230 and b >= 230


def is_gray(rgb: Optional[Tuple[int, int, int]]) -> bool:
    """灰色（非主题色）：三通道数值接近相等（差 ≤ 25），亮度不至于纯黑或纯白。"""
    if rgb is None:
        return False
    r, g, b = rgb
    if max(r, g, b) - min(r, g, b) > 25:
        return False
    return 40 <= min(r, g, b) <= 230


def is_light_color(rgb: Optional[Tuple[int, int, int]]) -> bool:
    """浅色（非主题色）：最小通道 ≥ 180，即整体明亮的低饱和色。"""
    if rgb is None:
        return False
    r, g, b = rgb
    return min(r, g, b) >= 180


def is_non_theme_light(rgb: Optional[Tuple[int, int, int]]) -> bool:
    """非主题色的浅/白/灰颜色：符合白色 / 浅色 / 灰色其中之一即可。

    用于"非主题色颜色设定"—— 只要颜色本身符合要求即通过，无需与主题蓝色相近。
    """
    return is_white(rgb) or is_gray(rgb) or is_light_color(rgb)


def shape_fill_rgb(shape) -> Optional[Tuple[int, int, int]]:
    """尝试取得形状填充色 RGB。"""
    try:
        fill = shape.fill
        if fill is None:
            return None
        # 渐变 / 图片 / 主题色 等情况 .fore_color.rgb 可能报错
        rgb = fill.fore_color.rgb
        if rgb is None:
            return None
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


def shape_has_gradient_fill(shape) -> bool:
    """看 XML 是否使用 gradFill。"""
    try:
        xml = shape._element.xml
        return "gradFill" in xml
    except Exception:
        return False


def shape_line_width_pt(shape) -> Optional[float]:
    try:
        w = shape.line.width
        if w is None:
            return None
        return emu_to_pt(w)
    except Exception:
        return None


def shape_line_rgb(shape) -> Optional[Tuple[int, int, int]]:
    try:
        rgb = shape.line.color.rgb
        if rgb is None:
            return None
        return (rgb[0], rgb[1], rgb[2])
    except Exception:
        return None


# ---------- 几何工具 ----------
def shape_bbox_cm(shape) -> Tuple[float, float, float, float]:
    """返回 (left, top, right, bottom)，单位 cm。"""
    l = emu_to_cm(shape.left or 0)
    t = emu_to_cm(shape.top or 0)
    w = emu_to_cm(shape.width or 0)
    h = emu_to_cm(shape.height or 0)
    return l, t, l + w, t + h


def rect_intersect_area(a, b) -> float:
    al, at, ar, ab = a
    bl, bt, br, bb = b
    iw = max(0.0, min(ar, br) - max(al, bl))
    ih = max(0.0, min(ab, bb) - max(at, bt))
    return iw * ih


def rect_area(a) -> float:
    return max(0.0, a[2] - a[0]) * max(0.0, a[3] - a[1])


# ---------- 维度 1 ----------
def check_dim1(pptx_path: str) -> Tuple[bool, List[str]]:
    """维度1：
       1) 交付文件为 .pptx 格式，文件可正常打开；
       2) 交付PPT包含16页幻灯片。
    返回 (是否通过, 失败原因列表)。
    """
    fails: List[str] = []

    # 1. 文件存在 & 后缀必须为 .pptx
    if not os.path.exists(pptx_path):
        return False, [f"文件不存在: {pptx_path}"]
    if not pptx_path.lower().endswith(".pptx"):
        fails.append("文件后缀不是 .pptx")

    # 2. 文件可正常打开：能作为 zip 解包，且能被 python-pptx 解析
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            _ = zf.namelist()
    except zipfile.BadZipFile:
        return False, ["文件不是合法的 pptx (zip) 包，无法打开"]

    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return False, [f"python-pptx 解析失败，文件无法正常打开: {e}"]

    # 3. 交付PPT包含16页幻灯片
    slide_count = len(prs.slides)
    if slide_count != 16:
        fails.append(f"交付PPT幻灯片数量为 {slide_count}，不等于 16 页")

    return (len(fails) == 0), fails


# ---------- 标题栏识别 ----------
# ---------- 维度 2 ----------
def evaluate_dim2(pptx_path: str, log: Dim2Log):
    prs = Presentation(pptx_path)
    slide_w_cm = emu_to_cm(prs.slide_width)
    slide_h_cm = emu_to_cm(prs.slide_height)
    slides = list(prs.slides)

    # 针对 2-15 页 (索引 1..14)，先抽取标题栏关键元素
    page_header_info = []  # 每页一个 dict
    for i in range(1, 15):
        # page_num = 2..15，用于校验"页码/章节编号"与实际页码对应
        page_header_info.append(
            extract_header_info(slides[i], slide_w_cm, slide_h_cm, page_num=i + 1)
        )

    # +5：第2-15页PPT顶部标题外侧均添加蓝色形状组合标题栏
    #    按细则拆分为以下并列要点（每一页都需要全部满足）：
    #      a) 范围：第 2-15 页（共 14 页）；
    #      b) 位置：标题栏位于 PPT "顶部"（顶端区域）；
    #      c) 位置：位于原页面标题的"外侧"（不遮挡/不覆盖原标题，在其上方/外部）；
    #      d) 形态：为"形状组合"——由多个形状组成（≥2 个形状）；
    #      e) 颜色：整体呈蓝色（包含蓝色形状元素）。
    cnt_top_outside_blue_group = sum(
        1 for info in page_header_info if info["top_outside_blue_group_ok"]
    )
    log.record(
        "第2-15页PPT顶部标题外侧均添加蓝色形状组合标题栏",
        5,
        cnt_top_outside_blue_group == 14,
        f"达标页数 {cnt_top_outside_blue_group}/14",
    )

    # +5: 左侧蓝色主形状（梯形/折角/平行四边形），宽3-4cm，高2-3cm，蓝色或浅蓝色（含渐变），轮廓可编辑
    ok_left_shape = sum(1 for info in page_header_info if info["left_main_ok"])
    log.record(
        "第2-15页标题栏左侧蓝色主形状（宽3-4cm、高2-3cm、蓝色/浅蓝/渐变填充）均达标",
        5,
        ok_left_shape == 14,
        f"达标页数 {ok_left_shape}/14",
    )

    # +5: 右侧长条外框 宽12-15cm 高1.5-3cm 蓝细线 + 斜切过渡 + 所有标题栏形状线宽 0.75-1.5pt 蓝色
    ok_right = sum(1 for info in page_header_info if info["right_bar_ok"])
    log.record(
        "第2-15页标题栏右侧长条外框（宽12-15cm、高1.5-3cm、蓝细实线）、斜切过渡形状与线宽均达标",
        5,
        ok_right == 14,
        f"达标页数 {ok_right}/14",
    )

    # +5: 文本元素 —— 页码/编号(Arial 28-32pt 白/浅色)、中文标题(宋体加粗 30pt 蓝色)、英文副标题(Arial 10-14pt 浅蓝)
    ok_text = sum(1 for info in page_header_info if info["text_ok"])
    log.record(
        "第2-15页标题栏文字（页码 Arial 28-32pt 白/浅色；中文宋体加粗 30pt 蓝色；英文 Arial 10-14pt 浅蓝）均达标",
        5,
        ok_text == 14,
        f"达标页数 {ok_text}/14",
    )

    # +5: 装饰点阵(9-25 圆点) + 3 条短斜线，位于右侧、不重叠、线宽1-3pt
    ok_decor = sum(1 for info in page_header_info if info["decor_ok"])
    log.record(
        "第2-15页标题栏右侧装饰点阵（9-25个圆点）与 3 条浅蓝色短斜线均达标",
        5,
        ok_decor == 14,
        f"达标页数 {ok_decor}/14",
    )

    # +1：第16页 "谢谢观看" 文本
    if len(slides) >= 16:
        end_ok, end_detail = check_end_thanks(slides[15], slide_w_cm, slide_h_cm)
    else:
        end_ok, end_detail = False, "不足16页"
    log.record(
        "第16页'谢谢观看'文本（宋体、加粗、36-50pt、蓝色、居中且位于上方区域）达标",
        1,
        end_ok,
        end_detail,
    )

    # +5：第16页视频对象（左右两段视频、尺寸、对齐、间距）
    if len(slides) >= 16:
        vid_ok, vid_detail = check_end_videos(slides[15], slide_w_cm, slide_h_cm)
    else:
        vid_ok, vid_detail = False, "不足16页"
    log.record(
        "第16页插入了左右两段视频，宽8-10cm、高11-12cm、顶部对齐、间距1.0-2.5cm、居中位于页面中下部",
        5,
        vid_ok,
        vid_detail,
    )

    # +5：视频可播放 + 静态封面/首帧
    media_ok, media_detail = check_end_videos_playable(pptx_path)
    log.record(
        "第16页两段视频在播放模式下可正常播放，且具备静态封面/首帧图像",
        5,
        media_ok,
        media_detail,
    )



# ---------- 标题栏内部判定 ----------
def extract_header_info(slide, slide_w_cm: float, slide_h_cm: float, page_num: int = 0) -> dict:
    """从单页中提取标题栏的各项检测结论。

    page_num：该页的实际页码（如 2, 3, ..., 15）；用于校验标题栏左侧
    "页码/章节编号" 是否与当前页对应。0 表示不校验页码对应关系（保留兼容）。
    """
    header_shapes = [sh for sh in slide.shapes if emu_to_cm(sh.top or 0) < 3.0]

    # ---------- top_outside_blue_group_ok：识别"标题栏形状集合"本身 ----------
    # 思路：不再把"顶部 20% 内任意形状"当作标题栏，而是先识别真正构成"标题栏"的
    # 形状集合 —— 位于页面顶部区域、具备蓝色特征（蓝色/浅蓝填充或描边、含蓝色渐变）
    # 的 AutoShape 集合（排除图表、媒体、以及横跨整页的大图/背景）。
    # 再从"标题栏形状集合本身"出发，验证：
    #   a) 位置：整体位于页面顶部；
    #   b) 蓝色组合：集合内至少有 2 个蓝色形状，形成"蓝色形状组合"；
    #   c) 形状组合：由多个形状组成（≥2 个形状）；
    #   d) 水平延展：集合的水平跨度足以构成"标题栏"（≥ 页宽 50%）；
    #   e) 外侧关系：原页面主标题不被标题栏形状集合覆盖 ——
    #      主标题（含中文、字号 ≥18pt 的文本框）与标题栏集合 bbox 的垂直重叠
    #      不超过主标题自身高度的 30%。
    top_zone_limit = slide_h_cm * 0.20

    def _shape_is_blueish(sh):
        fill_rgb = shape_fill_rgb(sh)
        line_rgb = shape_line_rgb(sh)
        return (
            is_blue(fill_rgb)
            or is_light_blue_or_white(fill_rgb)
            or is_blue(line_rgb)
            or is_light_blue_or_white(line_rgb)
            or shape_has_gradient_fill(sh)
        )

    title_bar_shapes = []
    for sh in slide.shapes:
        if sh.shape_type in (MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.MEDIA):
            continue
        top_cm = emu_to_cm(sh.top or 0)
        if top_cm >= top_zone_limit:
            continue
        w_cm = emu_to_cm(sh.width or 0)
        h_cm = emu_to_cm(sh.height or 0)
        # 排除横跨整页的大图片（可能是背景 / 装饰底图，不属于标题栏组合）
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE and w_cm >= slide_w_cm * 0.8:
            continue
        if _shape_is_blueish(sh):
            title_bar_shapes.append(sh)

    # c) 形状组合：≥2 个形状；b) 蓝色组合：集合本身即为蓝色形状（已在筛选时保证）
    group_ok = len(title_bar_shapes) >= 2
    has_blue_element = len(title_bar_shapes) >= 2

    # a) / d) 标题栏集合的相对位置与水平跨度
    if title_bar_shapes:
        tb_left = min(emu_to_cm(sh.left or 0) for sh in title_bar_shapes)
        tb_top = min(emu_to_cm(sh.top or 0) for sh in title_bar_shapes)
        tb_right = max(
            emu_to_cm((sh.left or 0) + (sh.width or 0)) for sh in title_bar_shapes
        )
        tb_bottom = max(
            emu_to_cm((sh.top or 0) + (sh.height or 0)) for sh in title_bar_shapes
        )
        span_ok = (tb_right - tb_left) >= slide_w_cm * 0.5
        # 集合整体位于顶部：bbox 顶端接近页面顶部（tb_top < 页高 15%）
        top_position_ok = tb_top < slide_h_cm * 0.15
    else:
        tb_left = tb_top = tb_right = tb_bottom = 0
        span_ok = False
        top_position_ok = False

    # e) 外侧关系：原页面主标题不被标题栏形状集合覆盖
    #    ——识别整页中"含中文、字号 ≥18pt"的文本框作为主标题候选，
    #    候选文本框不属于标题栏集合，其 bbox 与标题栏集合 bbox 的垂直重叠
    #    应不超过其自身高度的 30%。
    outside_ok = True
    title_bar_set = set(id(s) for s in title_bar_shapes)
    for sh in slide.shapes:
        if id(sh) in title_bar_set:
            continue
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if not txt:
            continue
        has_cn = any("一" <= ch <= "鿿" for ch in txt)
        if not has_cn:
            continue
        try:
            run = sh.text_frame.paragraphs[0].runs[0]
            size_pt = run.font.size.pt if run.font.size else 0
        except Exception:
            size_pt = 0
        if size_pt < 18:
            continue  # 非主标题级别字号
        l, t, r, b = shape_bbox_cm(sh)
        if b - t <= 0.01 or not title_bar_shapes:
            continue
        overlap_h = max(0.0, min(b, tb_bottom) - max(t, tb_top))
        if overlap_h > (b - t) * 0.3:
            outside_ok = False
            break

    top_outside_blue_group_ok = (
        group_ok
        and has_blue_element
        and top_position_ok
        and span_ok
        and outside_ok
    )

    # ---------- has_header：保留旧字段（供其它评分项内部参考） ----------
    blue_like_cnt = 0
    for sh in header_shapes:
        rgb = shape_fill_rgb(sh)
        if is_blue(rgb) or is_light_blue_or_white(rgb) or shape_has_gradient_fill(sh):
            blue_like_cnt += 1
    has_header = blue_like_cnt >= 3

    # ---------- left_main_ok：左侧蓝色主形状 ----------
    # 按细则拆分为以下并列要点（每一页都需要全部满足）：
    #   a) 位置：位于"标题栏左侧"；
    #   b) 形态：为"梯形"或"折角卡片形状"（折角卡片可视为平行四边形/带斜边的卡片）；
    #   c) 宽度：3-4 厘米；
    #   d) 高度：2-3 厘米；
    #   e) 填充：蓝色渐变 或 浅蓝色；
    #   f) 轮廓可编辑（即非锁定 / 可读取到 line 属性）；
    #   g) 唯一性：标题栏左侧"仅包含一个"满足上述形态(梯形/折角卡片)的形状（不能出现 2 个或以上）。
    import re
    card_shapes = {
        "trapezoid",
        "parallelogram",
        "foldedCorner",
        "flowChartManualOperation",
        "flowChartManualInput",
    }

    # g) 先统计标题栏左侧区域内"梯形/折角卡片"形状总数
    left_card_shapes_in_zone = []
    for sh in header_shapes:
        l, t, r, b = shape_bbox_cm(sh)
        if l > slide_w_cm * 0.35:  # a) 标题栏左侧
            continue
        xml = sh._element.xml
        m = re.search(r'prstGeom prst="([^"]+)"', xml)
        prst = m.group(1) if m else ""
        if prst in card_shapes:
            left_card_shapes_in_zone.append(sh)

    left_main_ok = False
    # 仅当左侧区域恰好有 1 个梯形/折角卡片形状时才继续后续判定
    if len(left_card_shapes_in_zone) == 1:
        sh = left_card_shapes_in_zone[0]
        l, t, r, b = shape_bbox_cm(sh)
        w = r - l
        h = b - t
        ok = True
        # c) 宽度 3-4cm
        if not (3.0 <= w <= 4.0):
            ok = False
        # d) 高度 2-3cm
        if ok and not (2.0 <= h <= 3.0):
            ok = False
        # e) 填充：蓝色渐变 或 浅蓝色
        if ok:
            has_grad = shape_has_gradient_fill(sh)
            rgb = shape_fill_rgb(sh)
            light_blue_hit = False
            if rgb is not None:
                rd, gn, bl = rgb
                # 浅蓝：B 通道较高且整体偏亮（R/G 也较大但不超过 B）
                if bl >= 180 and bl >= rd and bl >= gn - 20 and (rd + gn) >= 240:
                    light_blue_hit = True
            if not (has_grad or light_blue_hit):
                ok = False
        # f) 轮廓可编辑：能够访问 sh.line 即视为可编辑
        if ok:
            try:
                _ = sh.line
            except Exception:
                ok = False
        left_main_ok = ok

    # ---------- right_bar_ok：右侧长条外框 + 斜切过渡 + 标题栏整体线宽颜色 ----------
    # 按细则拆分为以下并列要点（每一页都需要全部满足）：
    #   A. 右侧长条外框（"长圆角平行四边形外框"）：
    #      A1) 位于标题栏"右侧"；
    #      A2) 形态：长圆角平行四边形外框；
    #      A3) 颜色：填充为浅蓝色 或 白色；
    #      A4) 宽度 12-15 cm；
    #      A5) 高度 1.5-3 cm；
    #      A6) 边框：蓝色细实线。
    #   B. 斜切过渡形状（位于左侧主形状与右侧长条之间）：
    #      B1) 位置：处于左侧主形状与右侧长条之间；
    #      B2) 角度约 15°-30°；
    #      B3) 颜色：白色 或 浅蓝色；
    #      B4) 形成层叠卡片效果（与左/右元素有水平交叠）。
    #   C. 标题栏整体形状的边框：
    #      C1) 线宽 0.75-1.5 磅；
    #      C2) 颜色：蓝色 或 浅蓝色。
    import re

    def _prst_of(sh):
        m = re.search(r'prstGeom prst="([^"]+)"', sh._element.xml)
        return m.group(1) if m else ""

    def _is_solid_blue_line(sh):
        # 边框颜色只接受"蓝色 或 浅蓝色"，不再接受白色。
        line_rgb = shape_line_rgb(sh)
        if not is_blue_or_light_blue(line_rgb):
            return False
        xml = sh._element.xml
        m_dash = re.search(r'prstDash\s+val="([^"]+)"', xml)
        if m_dash and m_dash.group(1) != "solid":
            return False
        return True

    def _is_white_or_light_blue_fill(sh):
        # "填充为浅蓝色 或 白色"：仅接受浅蓝色 或 白色，不再放宽到灰色/其它浅色。
        if shape_has_gradient_fill(sh):
            return False
        rgb = shape_fill_rgb(sh)
        if rgb is None:
            return False
        # 白色：三通道均 ≥230
        if is_white(rgb):
            return True
        # 浅蓝：不含白色，且 B 通道占优
        if is_light_blue(rgb):
            return True
        return False

    # A) 右侧长条外框
    right_bar_found = False
    right_bar_shape = None
    parallelogram_card_shapes = {
        "parallelogram",
        "roundRect",
        "round1Rect",
        "round2SameRect",
        "round2DiagRect",
        "snip1Rect",
        "snip2SameRect",
        "snip2DiagRect",
        "snipRoundRect",
    }
    for sh in header_shapes:
        l, t, r, b = shape_bbox_cm(sh)
        w = r - l
        h = b - t
        if (l + w / 2) < slide_w_cm * 0.5:  # A1) 右侧
            continue
        if not (12.0 <= w <= 15.0):  # A4)
            continue
        if not (1.5 <= h <= 3.0):  # A5)
            continue
        if _prst_of(sh) not in parallelogram_card_shapes:  # A2)
            continue
        if not _is_white_or_light_blue_fill(sh):  # A3)
            continue
        if not _is_solid_blue_line(sh):  # A6)
            continue
        right_bar_found = True
        right_bar_shape = sh
        break

    # B) 斜切过渡形状
    slant_found = False
    left_main_bbox = None
    for sh in header_shapes:
        l, t, r, b = shape_bbox_cm(sh)
        if l > slide_w_cm * 0.35:
            continue
        w = r - l
        h = b - t
        if not (3.0 <= w <= 4.0 and 2.0 <= h <= 3.0):
            continue
        if _prst_of(sh) not in {
            "trapezoid",
            "parallelogram",
            "foldedCorner",
            "flowChartManualOperation",
            "flowChartManualInput",
        }:
            continue
        left_main_bbox = (l, t, r, b)
        break

    for sh in header_shapes:
        if _prst_of(sh) != "parallelogram":
            continue
        l, t, r, b = shape_bbox_cm(sh)
        w = r - l
        h = b - t
        if w <= 0 or h <= 0:
            continue
        # B2) 角度：解析 parallelogram 的 adj 值（PPT 几何：tan(angle)≈adj）。
        #     15°→adj≈0.268；30°→adj≈0.577。给一点容差 0.20-0.62。
        xml = sh._element.xml
        m_adj = re.search(r'<a:gd\s+name="adj"\s+fmla="val\s+(-?\d+)"', xml)
        if m_adj:
            adj = int(m_adj.group(1)) / 100000.0
        else:
            adj = 0.25  # parallelogram 默认值
        if not (0.20 <= abs(adj) <= 0.62):
            continue
        if not _is_white_or_light_blue_fill(sh):  # B3)
            continue
        # B1) 位于左侧主形状与右侧长条之间
        if left_main_bbox is not None and right_bar_shape is not None:
            lm_right = left_main_bbox[2]
            rb_l = emu_to_cm(right_bar_shape.left or 0)
            rb_r = rb_l + emu_to_cm(right_bar_shape.width or 0)
            cx = (l + r) / 2
            if not (lm_right - 1.0 <= cx <= rb_r):
                continue
        # B4) 层叠卡片效果：与左侧主形状或右侧长条在水平上有重叠
        stacked = False
        if left_main_bbox is not None:
            if not (r < left_main_bbox[0] or l > left_main_bbox[2]):
                stacked = True
        if (not stacked) and right_bar_shape is not None:
            rb_l = emu_to_cm(right_bar_shape.left or 0)
            rb_r = rb_l + emu_to_cm(right_bar_shape.width or 0)
            if not (r < rb_l or l > rb_r):
                stacked = True
        if not stacked:
            continue
        slant_found = True
        break

    # C) 所有标题栏形状的可见边框必须显式满足：
    #    - 线宽 0.75-1.5 磅（须显式设置且 > 0）
    #    - 颜色为蓝色 或 浅蓝色（不接受白色/灰色）
    # 判定策略：先识别"可见边框"的标题栏形状 —— 描边线宽已显式设置且 > 0，
    # 或 XML 中定义了 <a:ln> 且未声明 noFill/noStroke。对这些形状逐一验证；
    # 若识别到 0 个可见边框（异常情况），亦判为不满足（无法证明"均达标"）。
    line_widths_ok = True
    visible_border_count = 0
    for sh in header_shapes:
        lw = shape_line_width_pt(sh)
        # 检测该形状是否显式声明了描边（<a:ln>...）且非 noFill
        try:
            sh_xml = sh._element.xml
        except Exception:
            sh_xml = ""
        has_ln_elem = "<a:ln" in sh_xml
        no_stroke = 'a:noFill' in sh_xml and "<a:ln" in sh_xml and (
            re.search(r'<a:ln\b[^>]*>\s*<a:noFill', sh_xml) is not None
        )
        has_visible_border = (lw is not None and lw > 0) or (has_ln_elem and not no_stroke)
        if not has_visible_border:
            continue
        visible_border_count += 1
        # 线宽必须显式在 0.75-1.5pt
        if lw is None or lw <= 0 or not (0.75 <= lw <= 1.5):
            line_widths_ok = False
            break
        # 颜色必须为蓝色 或 浅蓝色（不再接受白色）
        line_rgb = shape_line_rgb(sh)
        if not is_blue_or_light_blue(line_rgb):
            line_widths_ok = False
            break
    if visible_border_count == 0:
        line_widths_ok = False

    right_bar_ok = right_bar_found and slant_found and line_widths_ok

    # ---------- text_ok：页码 / 中文标题 / 英文副标题 ----------
    # 按细则拆分为以下并列要点（每一页都需要全部满足）：
    #   I. 页码 / 章节编号（如 02、03、04 等）
    #      I1) 显示对应页码/章节编号；
    #      I2) 字体 Arial；
    #      I3) 颜色为白色或浅色；
    #      I4) 字号 28-32 磅；
    #      I5) 位于"标题栏左侧区域"；
    #      I6) 位于"蓝色主形状内部"（被左侧蓝色主形状的边界框包含）。
    #   II. 中文页面标题
    #      II1) 内容含中文；
    #      II2) 字体宋体；
    #      II3) 加粗；
    #      II4) 颜色为蓝色；
    #      II5) 字号 30 磅；
    #      II6) 位于"长条中部偏左区域"（右侧长条 horizontal 范围内、偏左半侧）；
    #      II7) 文字不超出标题栏边界（文本框完全在 header_shapes 整体外接矩形内）。
    #   III. 英文副标题
    #      III1) 内容为英文（含英文字符）；
    #      III2) 字体 Arial；
    #      III3) 颜色为浅蓝色；
    #      III4) 字号 10-14 磅；
    #      III5) 位于中文标题下方或附近；
    #      III6) 文字内容与"页面主题"相对应（与同页中文标题相关）。

    # 先重新定位"左侧蓝色主形状"的 bbox 与"右侧长条"的 bbox，用于位置判定
    blue_main_bbox = None
    for sh in header_shapes:
        l, t, r, b = shape_bbox_cm(sh)
        if l > slide_w_cm * 0.35:
            continue
        w = r - l
        h = b - t
        if not (3.0 <= w <= 4.0 and 2.0 <= h <= 3.0):
            continue
        if _prst_of(sh) not in {
            "trapezoid",
            "parallelogram",
            "foldedCorner",
            "flowChartManualOperation",
            "flowChartManualInput",
        }:
            continue
        blue_main_bbox = (l, t, r, b)
        break

    long_bar_bbox = None
    for sh in header_shapes:
        l, t, r, b = shape_bbox_cm(sh)
        w = r - l
        h = b - t
        nm_l = sh.name.lower()
        if (
            ("rounded rectangle" in nm_l or "parallelogram" in nm_l or "rectangle" in nm_l)
            and w >= 12.0
            and 1.5 <= h <= 3.0
            and (l + w / 2) >= slide_w_cm * 0.4
        ):
            long_bar_bbox = (l, t, r, b)
            break

    # 标题栏整体外接矩形（用于"文字不超出标题栏边界"）
    if header_shapes:
        hb_l = min(emu_to_cm(sh.left or 0) for sh in header_shapes)
        hb_t = min(emu_to_cm(sh.top or 0) for sh in header_shapes)
        hb_r = max(
            emu_to_cm((sh.left or 0) + (sh.width or 0)) for sh in header_shapes
        )
        hb_b = max(
            emu_to_cm((sh.top or 0) + (sh.height or 0)) for sh in header_shapes
        )
    else:
        hb_l = hb_t = hb_r = hb_b = 0

    page_num_ok = False
    cn_title_ok = False
    en_sub_ok = False
    cn_title_text = ""       # 用于英文副标题主题对应判断
    cn_title_bbox = None     # 用于英文副标题"位于中文标题下方或附近"判定

    # 预扫描：先定位中文标题的位置和文本，再判定英文副标题时依赖它。
    # 记录 (shape, txt, size_pt, font_name, rgb, bold, bbox)
    text_items = []
    for sh in header_shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if not txt:
            continue
        runs = []
        for para in sh.text_frame.paragraphs:
            runs.extend(para.runs)
        if not runs:
            continue
        run = runs[0]
        font_name = (run.font.name or "")
        size_pt = run.font.size.pt if run.font.size else None
        rgb = _rgb_of(run.font)
        bold = run.font.bold
        bbox = shape_bbox_cm(sh)
        text_items.append((sh, txt, size_pt, font_name, rgb, bold, bbox))

    # 页码 / 章节编号的目标字符串（当传入 page_num≥2 时才严格校验对应）
    def _matches_page_num_text(txt: str, page_num: int) -> bool:
        if page_num <= 0:
            # 未指定实际页码 —— 退化为原来的宽松判定
            return (txt.isdigit() and 1 <= len(txt) <= 3) or (
                len(txt) <= 6
                and any(ch.isdigit() for ch in txt)
                and ("章" in txt or "节" in txt or txt.isdigit())
            )
        # 严格：显式对应当前页码
        stripped = txt.strip()
        expected_2 = f"{page_num:02d}"  # "02" .. "15"
        expected_1 = str(page_num)      # "2" .. "15"
        if stripped in (expected_1, expected_2):
            return True
        # 章节编号：形如 "第2章"、"第02章"、"第2节"
        cn_num_map = {
            2: "二", 3: "三", 4: "四", 5: "五", 6: "六", 7: "七",
            8: "八", 9: "九", 10: "十", 11: "十一", 12: "十二",
            13: "十三", 14: "十四", 15: "十五",
        }
        cn_num = cn_num_map.get(page_num, "")
        chapter_patterns = [
            f"第{expected_1}章", f"第{expected_2}章",
            f"第{expected_1}节", f"第{expected_2}节",
        ]
        if cn_num:
            chapter_patterns.extend([f"第{cn_num}章", f"第{cn_num}节"])
        for pat in chapter_patterns:
            if pat in stripped:
                return True
        return False

    # 第一轮：先扫描"页码 / 章节编号"和"中文标题"，把中文标题的 bbox / 文本记下来
    for sh, txt, size_pt, font_name, rgb, bold, (sl, st, sr, sb) in text_items:
        scx = (sl + sr) / 2
        scy = (st + sb) / 2

        is_page_num_text = _matches_page_num_text(txt, page_num)
        if is_page_num_text:
            cond_font = "arial" in font_name.lower()          # I2)
            cond_color = is_non_theme_light(rgb)              # I3)
            cond_size = size_pt is not None and 28 <= size_pt <= 32  # I4)
            cond_left = scx < slide_w_cm * 0.35               # I5)
            cond_inside_blue = False                          # I6)
            if blue_main_bbox is not None:
                bl, bt, br, bb = blue_main_bbox
                if bl <= scx <= br and bt <= scy <= bb:
                    cond_inside_blue = True
            if cond_font and cond_color and cond_size and cond_left and cond_inside_blue:
                page_num_ok = True

        has_cn = any("一" <= ch <= "鿿" for ch in txt)
        if has_cn and not is_page_num_text:
            cond_font = (
                "宋" in font_name
                or "Song" in font_name
                or "SimSun" in font_name
                or font_name == ""
            )                                                 # II2)
            cond_bold = bool(bold)                            # II3)
            cond_color = is_blue(rgb)                         # II4)
            cond_size = size_pt is not None and abs(size_pt - 30) < 0.5  # II5)
            cond_pos = True                                   # II6)
            if long_bar_bbox is not None:
                bl, bt, br, bb = long_bar_bbox
                bar_mid = (bl + br) / 2
                if not (bl <= scx <= bar_mid + 1.0 and bt <= scy <= bb):
                    cond_pos = False
            cond_inside = (                                   # II7)
                sl >= hb_l - 0.1
                and sr <= hb_r + 0.1
                and st >= hb_t - 0.1
                and sb <= hb_b + 0.1
            )
            if cond_font and cond_bold and cond_color and cond_size and cond_pos and cond_inside:
                cn_title_ok = True
                cn_title_text = txt
                cn_title_bbox = (sl, st, sr, sb)

    # 中文标题主题对应关键词（同时用于中文标题内容校验 与 英文副标题主题匹配）
    theme_keywords_cn = [
        "咖啡", "晨岚", "品牌", "运营", "研究", "市场", "分析", "调研",
        "定位", "策略", "产品", "推广", "营销", "价格", "渠道", "消费者",
        "目标", "总结", "背景", "现状", "规划",
    ]
    theme_keywords_en = [
        "coffee", "brand", "research", "operation", "chenlan",
        "café", "study", "analysis", "market",
    ]

    # 中文标题"内容与页面主题对应"—— 至少命中主题相关关键词 或 与英文副标题共享主题词
    if cn_title_ok and cn_title_text:
        if not any(k in cn_title_text for k in theme_keywords_cn):
            cn_title_ok = False

    # 第二轮：英文副标题
    for sh, txt, size_pt, font_name, rgb, bold, (sl, st, sr, sb) in text_items:
        scx = (sl + sr) / 2
        scy = (st + sb) / 2
        has_cn = any("一" <= ch <= "鿿" for ch in txt)
        if has_cn:
            continue
        has_en = any("a" <= ch.lower() <= "z" for ch in txt)
        if not has_en:
            continue
        cond_font = "arial" in font_name.lower()              # III2)
        cond_color = is_light_blue(rgb)                       # III3) 浅蓝色（排除白色）
        cond_size = size_pt is not None and 10 <= size_pt <= 14  # III4)

        # III5) 位于中文标题下方或附近：
        #   要求英文副标题文本框的顶部不高于中文标题的中线，
        #   且水平方向与中文标题存在交叠或距离 ≤2cm（"附近"）。
        cond_pos = False
        if cn_title_ok and cn_title_bbox is not None:
            cl, ct, cr, cb = cn_title_bbox
            cn_mid_y = (ct + cb) / 2
            # 下方：副标题顶部 st 大于中文标题中线（允许 0.2cm 容差）
            below_ok = st >= cn_mid_y - 0.2
            # 附近：水平方向有交叠 或 中心距 ≤2cm
            horiz_overlap = not (sr < cl or sl > cr)
            horiz_near = abs(scx - (cl + cr) / 2) <= 2.0
            cond_pos = below_ok and (horiz_overlap or horiz_near)

        # III6) 主题对应：与中文标题共享主题（关键词命中）
        cond_topic = any(k in txt.lower() for k in theme_keywords_en)

        if cond_font and cond_color and cond_size and cond_pos and cond_topic:
            en_sub_ok = True
            break

    text_ok = page_num_ok and cn_title_ok and en_sub_ok

    # ---------- decor_ok：点阵装饰 + 3 条短斜线装饰 ----------
    # 按细则拆分为以下并列要点（每一页都需要全部满足）：
    #   D. 浅蓝色点阵装饰：
    #      D1) 包含点阵装饰；
    #      D2) 颜色为浅蓝色；
    #      D3) 由 9-25 个小圆点组成；
    #      D4) 位于"标题栏右上区域"（标题栏右半侧 + 偏上 1/2）；
    #      D5) 圆点大小一致。
    #   E. 短斜线装饰：
    #      E1) 共 3 条；
    #      E2) 颜色为浅蓝色；
    #      E3) 位于"标题栏右下方"（标题栏右半侧 + 偏下 1/2）；
    #      E4) 与点阵装饰不重叠；
    #      E5) 线宽 1-3 磅（包含线宽 或 形状窄边的视觉宽度）；
    #      E6) 角度一致；
    #      E7) 不遮挡正文（不与标题栏外的正文元素相交）。

    # 标题栏整体外接矩形（右上 / 右下区域划分依据）
    # 注：header_shapes 是按 top<3cm 筛出的；但有些页面装饰图片紧贴标题栏顶部（top≈2.5cm），
    # 会把 bbox 的 bottom 拉得过低。这里只用"真正属于标题装饰带"的形状来界定 bbox ——
    # 即排除大图片/正文图表，只保留 AutoShape（线段、圆、卡片等）。
    decor_band = [
        sh
        for sh in header_shapes
        if sh.shape_type not in (
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.CHART,
            MSO_SHAPE_TYPE.MEDIA,
        )
    ]
    if decor_band:
        h_left = min(emu_to_cm(sh.left or 0) for sh in decor_band)
        h_top = min(emu_to_cm(sh.top or 0) for sh in decor_band)
        h_right = max(
            emu_to_cm((sh.left or 0) + (sh.width or 0)) for sh in decor_band
        )
        h_bottom = max(
            emu_to_cm((sh.top or 0) + (sh.height or 0)) for sh in decor_band
        )
    else:
        h_left = h_top = h_right = h_bottom = 0
    h_mid_x = (h_left + h_right) / 2
    h_mid_y = (h_top + h_bottom) / 2

    # 收集圆点与短斜线 —— 按 prstGeom / 线段类型识别，而不是仅看 shape.name。
    # 圆点：prstGeom = "ellipse" / "oval" 的 AutoShape，且宽高相近（≈圆形），尺寸小；
    # 短斜线：Connector（直线连接符）或 prstGeom = "line" / "straightConnector1"，
    #        或形态为窄长的 parallelogram（细长条 + 有斜切量），且尺寸小。
    def _prst(sh) -> str:
        try:
            m = re.search(r'prstGeom prst="([^"]+)"', sh._element.xml)
            return m.group(1) if m else ""
        except Exception:
            return ""

    def _is_line_connector(sh) -> bool:
        try:
            xml = sh._element.xml
        except Exception:
            return False
        return ("<p:cxnSp" in xml) or (
            "<a:ln " in xml and "prstGeom prst=\"line\"" in xml
        )

    dot_shapes = []
    slash_shapes = []
    for sh in header_shapes:
        l, t, r, b = shape_bbox_cm(sh)
        w = r - l
        h = b - t
        prst = _prst(sh)
        # ---- 圆点识别 ----
        # 依据 prstGeom = ellipse/oval；要求近似圆形（|w-h| < 0.15cm）；尺寸小（≤0.5cm）。
        is_dot_prst = prst in ("ellipse", "oval")
        if is_dot_prst and w <= 0.5 and h <= 0.5 and abs(w - h) <= 0.15 and w > 0 and h > 0:
            dot_shapes.append(sh)
            continue
        # ---- 短斜线识别 ----
        # (a) 直线连接符/直线：<p:cxnSp> 或 prstGeom = "line" / "straightConnector1"；
        # (b) 或者：窄长 parallelogram（细长条：min(w,h)<=0.3cm 且 max(w,h)<=1.2cm）。
        line_like = (
            _is_line_connector(sh)
            or prst in ("line", "straightConnector1")
        )
        narrow_para = (
            prst == "parallelogram"
            and min(w, h) <= 0.3
            and max(w, h) <= 1.2
        )
        if (line_like and max(w, h) <= 1.5 and min(w, h) <= 0.4) or narrow_para:
            slash_shapes.append(sh)

    # D2) 点阵颜色：浅蓝色
    dots_color_ok = all(
        _is_white_or_light_blue_fill(sh) and not _is_white_fill(sh)
        if False
        else (
            (shape_fill_rgb(sh) is not None)
            and (shape_fill_rgb(sh)[2] >= 180)
            and (shape_fill_rgb(sh)[2] >= shape_fill_rgb(sh)[0])
        )
        for sh in dot_shapes
    ) if dot_shapes else False
    # D3) 9-25 个
    dots_count_ok = 9 <= len(dot_shapes) <= 25
    # D4) 位于"标题栏右上区域"：所有点中心 X 偏右半侧；且点阵整体偏上（点阵平均 Y < 斜线平均 Y）
    dots_x_right = all(
        ((emu_to_cm(sh.left or 0) + emu_to_cm(sh.width or 0) / 2) >= h_mid_x)
        for sh in dot_shapes
    ) if dot_shapes else False
    if dot_shapes and slash_shapes:
        dots_avg_y = sum(
            (emu_to_cm(sh.top or 0) + emu_to_cm(sh.height or 0) / 2)
            for sh in dot_shapes
        ) / len(dot_shapes)
        slashes_avg_y = sum(
            (emu_to_cm(sh.top or 0) + emu_to_cm(sh.height or 0) / 2)
            for sh in slash_shapes
        ) / len(slash_shapes)
        dots_upper_relative = dots_avg_y <= slashes_avg_y
    else:
        dots_upper_relative = False
    dots_position_ok = dots_x_right and dots_upper_relative
    # D5) 圆点大小一致：所有点宽度方差小于 0.02cm，高度同
    if dot_shapes:
        ws = [emu_to_cm(sh.width or 0) for sh in dot_shapes]
        hs = [emu_to_cm(sh.height or 0) for sh in dot_shapes]
        dots_size_uniform = (max(ws) - min(ws) <= 0.05) and (max(hs) - min(hs) <= 0.05)
    else:
        dots_size_uniform = False
    dots_ok = dots_count_ok and dots_color_ok and dots_position_ok and dots_size_uniform

    # E1) 3 条
    slashes_count_ok = len(slash_shapes) == 3
    # E2) 颜色浅蓝色
    slashes_color_ok = all(
        (shape_fill_rgb(sh) is not None)
        and (shape_fill_rgb(sh)[2] >= 180)
        and (shape_fill_rgb(sh)[2] >= shape_fill_rgb(sh)[0])
        for sh in slash_shapes
    ) if slash_shapes else False
    # E3) 位于标题栏右下方
    slashes_position_ok = all(
        ((emu_to_cm(sh.left or 0) + emu_to_cm(sh.width or 0) / 2) >= h_mid_x)
        and ((emu_to_cm(sh.top or 0) + emu_to_cm(sh.height or 0) / 2) >= h_mid_y)
        for sh in slash_shapes
    ) if slash_shapes else False
    # E4) 与点阵不重叠：以"斜线中心是否落入圆点 bbox 内 或 圆点中心是否落入斜线 bbox 内"作为
    #     视觉重叠判据（斜线/圆点的 bounding box 矩形会包含空白边界，单纯 bbox 相交容易误判）。
    if dot_shapes and slash_shapes:
        no_overlap = True
        for sl in slash_shapes:
            sb = shape_bbox_cm(sl)
            scx = (sb[0] + sb[2]) / 2
            scy = (sb[1] + sb[3]) / 2
            for d in dot_shapes:
                db = shape_bbox_cm(d)
                dcx = (db[0] + db[2]) / 2
                dcy = (db[1] + db[3]) / 2
                # 斜线中心落入点 bbox
                if db[0] <= scx <= db[2] and db[1] <= scy <= db[3]:
                    no_overlap = False
                    break
                # 点中心落入斜线 bbox
                if sb[0] <= dcx <= sb[2] and sb[1] <= dcy <= sb[3]:
                    no_overlap = False
                    break
            if not no_overlap:
                break
        slashes_no_overlap = no_overlap
    else:
        slashes_no_overlap = False
    # E5) 线宽严格 1-3 磅：
    #     - 直线连接符 / prstGeom=line 的斜线：使用 line.width（EMU→pt）判定；
    #     - 窄长 parallelogram 斜线：其视觉粗细 = 形状最短边，
    #       将其从 cm 换算为 pt（1cm≈28.35pt），必须落在 1-3pt 内。
    #     两种情况均严格限定 1-3pt，不再放宽到 12pt。
    slashes_lw_ok = True
    for sh in slash_shapes:
        prst = _prst(sh)
        is_line_shape = _is_line_connector(sh) or prst in ("line", "straightConnector1")
        lw = shape_line_width_pt(sh)
        if is_line_shape:
            # 直线：必须显式设置线宽，且落在 1-3pt
            if lw is None or lw <= 0 or not (1.0 <= lw <= 3.0):
                slashes_lw_ok = False
                break
        else:
            # 窄长 parallelogram 斜线：优先看描边线宽；若无描边则使用视觉短边
            if lw is not None and lw > 0:
                if not (1.0 <= lw <= 3.0):
                    slashes_lw_ok = False
                    break
            else:
                h_cm = emu_to_cm(sh.height or 0)
                w_cm = emu_to_cm(sh.width or 0)
                short_edge_pt = min(h_cm, w_cm) * 28.3464567
                if not (1.0 <= short_edge_pt <= 3.0):
                    slashes_lw_ok = False
                    break
    # E6) 角度一致：解析每条斜线的 adj 值（parallelogram 的斜切量），方差 ≤ 0.05；
    #     若无 adj（默认 25%），则要求形状宽高比一致（≤0.05 差）
    slashes_angle_ok = True
    if slash_shapes:
        adj_vals = []
        ratio_vals = []
        for sh in slash_shapes:
            xml = sh._element.xml
            m_adj = re.search(r'<a:gd\s+name="adj"\s+fmla="val\s+(-?\d+)"', xml)
            if m_adj:
                adj_vals.append(int(m_adj.group(1)) / 100000.0)
            w_cm = emu_to_cm(sh.width or 0)
            h_cm = emu_to_cm(sh.height or 0)
            ratio_vals.append(h_cm / w_cm if w_cm > 0 else 0)
        if adj_vals and len(adj_vals) == len(slash_shapes):
            if max(adj_vals) - min(adj_vals) > 0.05:
                slashes_angle_ok = False
        else:
            if max(ratio_vals) - min(ratio_vals) > 0.05:
                slashes_angle_ok = False
    else:
        slashes_angle_ok = False
    # E7) 不遮挡正文：与 top ≥ 标题栏底边 的非标题栏对象不相交
    slashes_no_block = True
    body_shapes_local = [
        sh
        for sh in slide.shapes
        if emu_to_cm(sh.top or 0) >= h_bottom - 0.1
    ]
    for sh in slash_shapes:
        sb = shape_bbox_cm(sh)
        for bsh in body_shapes_local:
            bb = shape_bbox_cm(bsh)
            iw = max(0.0, min(sb[2], bb[2]) - max(sb[0], bb[0]))
            ih = max(0.0, min(sb[3], bb[3]) - max(sb[1], bb[1]))
            if iw > 0.05 and ih > 0.05:
                slashes_no_block = False
                break
        if not slashes_no_block:
            break

    slashes_ok_overall = (
        slashes_count_ok
        and slashes_color_ok
        and slashes_position_ok
        and slashes_no_overlap
        and slashes_lw_ok
        and slashes_angle_ok
        and slashes_no_block
    )

    decor_ok = dots_ok and slashes_ok_overall

    return {
        "has_header": has_header,
        "top_outside_blue_group_ok": top_outside_blue_group_ok,
        "left_main_ok": left_main_ok,
        "right_bar_ok": right_bar_ok,
        "text_ok": text_ok,
        "decor_ok": decor_ok,
    }


# ---------- 第16页判定 ----------
def check_end_thanks(slide, slide_w_cm, slide_h_cm) -> Tuple[bool, str]:
    """第16页"谢谢观看"文本细则（每一点都必须满足）：
       1) 显示"谢谢观看"文本；
       2) 位于"页面上方或中上部"；
       3) 居中显示；
       4) 字体宋体；
       5) 字号 36-50 磅；
       6) 加粗；
       7) 颜色为蓝色；
       8) 文本框位于页面宽度 30%-70% 范围内；
       9) 文本框位于页面高度 10%-25% 范围内；
       10) 水平居中；
       11) 文字没有倾斜或变形（无旋转、无 X/Y 形变）。
    """
    # 1) 找到"谢谢观看"文本
    target = None
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if "谢谢观看" in sh.text_frame.text:
            target = sh
            break
    if target is None:
        return False, "未找到'谢谢观看'文本"

    runs = []
    for p in target.text_frame.paragraphs:
        runs.extend(p.runs)
    if not runs:
        return False, "'谢谢观看'文本框为空"
    run = runs[0]
    font_name = (run.font.name or "")
    size_pt = run.font.size.pt if run.font.size else None
    rgb = _rgb_of(run.font)
    bold = run.font.bold

    tl = emu_to_cm(target.left or 0)
    tt = emu_to_cm(target.top or 0)
    tw = emu_to_cm(target.width or 0)
    th = emu_to_cm(target.height or 0)
    cx_cm = tl + tw / 2
    cy_cm = tt + th / 2

    reasons = []

    # 4) 字体宋体
    if not (
        "宋" in font_name
        or "Song" in font_name
        or "SimSun" in font_name
        or font_name == ""
    ):
        reasons.append(f"字体非宋体({font_name})")
    # 5) 字号 36-50pt
    if size_pt is None or not (36 <= size_pt <= 50):
        reasons.append(f"字号 {size_pt} 不在 36-50pt")
    # 6) 加粗
    if not bold:
        reasons.append("未加粗")
    # 7) 颜色蓝色
    if not is_blue(rgb):
        reasons.append(f"颜色非蓝色 {rgb}")

    # 2) 位于"页面上方或中上部"：文本框中心 Y < 页面高度 50%
    if cy_cm > slide_h_cm * 0.5:
        reasons.append(
            f"文本未位于页面上方或中上部（中心 Y={cy_cm:.2f}cm > 页高 50%）"
        )

    # 3) 居中显示 + 10) 水平居中：文本框水平中线接近页面中线 (±1cm)
    if abs(cx_cm - slide_w_cm / 2) > 1.0:
        reasons.append(f"未水平居中（中心 X={cx_cm:.2f}cm 偏离页面中线>1cm）")
    # 段落对齐方式 = CENTER 进一步确认"居中显示"
    try:
        from pptx.enum.text import PP_ALIGN

        para_align_ok = any(
            p.alignment == PP_ALIGN.CENTER for p in target.text_frame.paragraphs
        )
        if not para_align_ok:
            # 文本框宽度若几乎横跨页面，也认可为视觉居中
            if tw < slide_w_cm * 0.6:
                reasons.append("段落对齐方式非居中")
    except Exception:
        pass

    # 8) 文本框位于页面宽度 30%-70% 范围内（文本框中心点）
    if not (slide_w_cm * 0.30 <= cx_cm <= slide_w_cm * 0.70):
        reasons.append(f"水平位置 {cx_cm:.2f}cm 不在页面 30%-70%")
    # 9) 文本框位于页面高度 10%-25% 范围内（文本框中心点）
    if not (slide_h_cm * 0.10 <= cy_cm <= slide_h_cm * 0.25):
        reasons.append(f"垂直位置 {cy_cm:.2f}cm 不在页面 10%-25%")

    # 11) 文字没有倾斜或变形（无旋转 & 无字符 X/Y 形变）
    try:
        if target.rotation and abs(target.rotation) > 0.5:
            reasons.append(f"文本框存在旋转 {target.rotation}°")
    except Exception:
        pass
    # 字符级倾斜/变形检查：italic 视为倾斜；XML 中 sp3d/scene3d/warp 视为变形
    if run.font.italic:
        reasons.append("文字被设置为倾斜(italic)")
    try:
        xml = target._element.xml
        if any(
            tag in xml
            for tag in ("<a:sp3d", "<a:scene3d", 'prst="textArchUp"', "prstTxWarp")
        ):
            reasons.append("文本存在 3D/变形(prstTxWarp)效果")
    except Exception:
        pass

    return (len(reasons) == 0), ("；".join(reasons) if reasons else "")


def _is_media(shape) -> bool:
    return shape.shape_type == MSO_SHAPE_TYPE.MEDIA


def check_end_videos(slide, slide_w_cm, slide_h_cm) -> Tuple[bool, str]:
    """第16页视频对象细则（每一点都必须满足）：
        L1) 左侧插入"咖啡制作"视频；
        L2) 位于"谢谢观看"下方左侧；
        L3) 左视频宽度 8-10cm；
        L4) 左视频高度 11-12cm；
        R1) 右侧插入"咖啡店环境介绍"视频；
        R2) 位于"谢谢观看"下方右侧；
        R3) 右视频宽度 8-10cm；
        R4) 右视频高度 11-12cm；
        P1) 两段视频左右并排；
        P2) 顶部对齐；
        P3) 高度一致或接近一致；
        P4) 水平间距 1.0-2.5cm；
        P5) 整体位于"页面中下部"；
        P6) 整体"居中"区域。
    """
    medias = [sh for sh in slide.shapes if _is_media(sh)]
    if len(medias) < 2:
        return False, f"页面媒体对象数量 {len(medias)} <2"
    medias = sorted(medias, key=lambda s: s.left or 0)
    left_v, right_v = medias[0], medias[1]

    reasons = []

    # 找到"谢谢观看"文本框，用于"下方左/右"判定
    thanks_box = None
    for sh in slide.shapes:
        if sh.has_text_frame and "谢谢观看" in sh.text_frame.text:
            thanks_box = sh
            break

    def _media_topic(media_shape) -> str:
        """从媒体对象提取主题名：先看 shape.name；再回退到关联媒体文件的 rId 名称。"""
        name = (media_shape.name or "").lower()
        # 关联文件名（如 'coffee_grind_to_latte_xxx.mp4'）
        try:
            xml = media_shape._element.xml
            # 形如 r:embed 或 r:link
            import re

            for m in re.finditer(r'r:(?:embed|link)="([^"]+)"', xml):
                rid = m.group(1)
                try:
                    part = media_shape.part.related_part(rid)
                    rel_name = (
                        getattr(part, "partname", "") or ""
                    ).lower() + " " + (
                        getattr(part, "filename", "") or ""
                    ).lower()
                    name += " " + rel_name
                except Exception:
                    pass
        except Exception:
            pass
        return name

    left_name = _media_topic(left_v)
    right_name = _media_topic(right_v)

    # L1) 左侧 = "咖啡制作"视频
    coffee_make_kw = [
        "咖啡制作",
        "制作",
        "grind",
        "latte",
        "brew",
        "coffee_grind",
        "coffee make",
        "coffee_making",
    ]
    if not any(k in left_name for k in coffee_make_kw):
        reasons.append(f"左侧视频不是'咖啡制作'相关 (name={left_name})")
    # R1) 右侧 = "咖啡店环境介绍"视频
    env_intro_kw = [
        "环境",
        "店",
        "咖啡店",
        "introduction",
        "intro",
        "shop",
        "环境介绍",
        "ambient",
    ]
    if not any(k in right_name for k in env_intro_kw):
        reasons.append(f"右侧视频不是'咖啡店环境介绍'相关 (name={right_name})")

    # L3 / R3 / L4 / R4：尺寸
    for label, v in (("左视频", left_v), ("右视频", right_v)):
        w = emu_to_cm(v.width or 0)
        h = emu_to_cm(v.height or 0)
        if not (8.0 <= w <= 10.0):
            reasons.append(f"{label}宽 {w:.2f}cm 不在 8-10cm")
        if not (11.0 <= h <= 12.0):
            reasons.append(f"{label}高 {h:.2f}cm 不在 11-12cm")

    # L2 / R2：位于"谢谢观看"下方左/右
    if thanks_box is not None:
        thanks_cx = emu_to_cm(thanks_box.left or 0) + emu_to_cm(thanks_box.width or 0) / 2
        thanks_bottom = emu_to_cm(thanks_box.top or 0) + emu_to_cm(thanks_box.height or 0)
        left_cx = emu_to_cm(left_v.left or 0) + emu_to_cm(left_v.width or 0) / 2
        right_cx = emu_to_cm(right_v.left or 0) + emu_to_cm(right_v.width or 0) / 2
        left_top = emu_to_cm(left_v.top or 0)
        right_top = emu_to_cm(right_v.top or 0)
        if not (left_cx < thanks_cx and left_top >= thanks_bottom - 0.2):
            reasons.append("左视频未位于'谢谢观看'下方左侧")
        if not (right_cx > thanks_cx and right_top >= thanks_bottom - 0.2):
            reasons.append("右视频未位于'谢谢观看'下方右侧")
    else:
        reasons.append("未找到'谢谢观看'文本，无法判断视频位置")

    # P1) 左右并排：左视频右边界 <= 右视频左边界（中间允许间距）
    left_right_edge = emu_to_cm(left_v.left or 0) + emu_to_cm(left_v.width or 0)
    right_left_edge = emu_to_cm(right_v.left or 0)
    if not (left_right_edge <= right_left_edge):
        reasons.append("两视频未呈左右并排")

    # P2) 顶部对齐（允许 ≤0.3cm 误差）
    top_diff = abs(emu_to_cm(left_v.top or 0) - emu_to_cm(right_v.top or 0))
    if top_diff > 0.3:
        reasons.append(f"两视频顶部差 {top_diff:.2f}cm，未对齐")

    # P3) 高度一致或接近一致（允许 ≤0.3cm）
    h_diff = abs(emu_to_cm(left_v.height or 0) - emu_to_cm(right_v.height or 0))
    if h_diff > 0.3:
        reasons.append(f"两视频高度差 {h_diff:.2f}cm，不一致")

    # P4) 水平间距 1.0-2.5cm
    gap = right_left_edge - left_right_edge
    if not (1.0 <= gap <= 2.5):
        reasons.append(f"两视频水平间距 {gap:.2f}cm 不在 1.0-2.5cm")

    # P5) 整体位于"页面中下部"：两视频上下垂直范围位于页面下半部分（顶 >= 50% 偏松到 30%）
    top_cm = min(emu_to_cm(left_v.top or 0), emu_to_cm(right_v.top or 0))
    bottom_cm = max(
        emu_to_cm(left_v.top or 0) + emu_to_cm(left_v.height or 0),
        emu_to_cm(right_v.top or 0) + emu_to_cm(right_v.height or 0),
    )
    mid_y = (top_cm + bottom_cm) / 2
    if not (mid_y >= slide_h_cm * 0.5):
        reasons.append(f"视频整体未位于页面中下部（中线 Y={mid_y:.2f}cm）")

    # P6) 整体"居中"：两视频整体水平中心接近页面中线 (≤1.5cm)
    overall_left = emu_to_cm(left_v.left or 0)
    overall_right = emu_to_cm(right_v.left or 0) + emu_to_cm(right_v.width or 0)
    overall_cx = (overall_left + overall_right) / 2
    if abs(overall_cx - slide_w_cm / 2) > 1.5:
        reasons.append(f"两视频整体中心 {overall_cx:.2f}cm 偏离页面中线")

    return (len(reasons) == 0), ("；".join(reasons) if reasons else "")


def _decode_first_frame_playable(video_bytes: bytes, ext: str) -> Tuple[bool, str]:
    """尝试用真实媒体解码器解码视频首帧，用于证明"可播放"。

    优先级：
      1) opencv-python (cv2)；
      2) 系统 ffprobe / ffmpeg 命令行（若可用）。
    任一路径能读出首帧的合法帧数据即视为"可播放"；全部失败则不通过。
    """
    import tempfile
    import subprocess
    import shutil

    # 写入临时文件（大多数解码器不接受内存流）
    tmp_dir = tempfile.mkdtemp(prefix="pptx_video_")
    tmp_path = os.path.join(tmp_dir, f"clip{ext}")
    try:
        with open(tmp_path, "wb") as f:
            f.write(video_bytes)

        # -------- 1) OpenCV --------
        try:
            import cv2  # type: ignore
            cap = cv2.VideoCapture(tmp_path)
            try:
                if cap.isOpened():
                    ok, frame = cap.read()
                    if ok and frame is not None and frame.size > 0:
                        return True, "cv2 解码首帧成功"
            finally:
                cap.release()
        except Exception:
            pass

        # -------- 2) ffprobe / ffmpeg CLI --------
        ffprobe = shutil.which("ffprobe")
        if ffprobe:
            try:
                cp = subprocess.run(
                    [
                        ffprobe, "-v", "error",
                        "-select_streams", "v:0",
                        "-show_entries", "stream=codec_name,width,height,duration",
                        "-of", "default=noprint_wrappers=1",
                        tmp_path,
                    ],
                    capture_output=True, text=True, timeout=15,
                )
                out = (cp.stdout or "") + (cp.stderr or "")
                if cp.returncode == 0 and "codec_name=" in out and "width=" in out:
                    return True, f"ffprobe 识别视频流成功：{out.strip()[:80]}"
            except Exception:
                pass
        ffmpeg = shutil.which("ffmpeg")
        if ffmpeg:
            try:
                # 用 ffmpeg 抽取首帧到 /dev/null 或临时文件
                out_frame = os.path.join(tmp_dir, "frame0.png")
                cp = subprocess.run(
                    [ffmpeg, "-y", "-i", tmp_path, "-vframes", "1", out_frame],
                    capture_output=True, timeout=20,
                )
                if cp.returncode == 0 and os.path.exists(out_frame) and os.path.getsize(out_frame) > 0:
                    return True, "ffmpeg 抽取首帧成功"
            except Exception:
                pass

        return False, "无可用媒体解码器（opencv-python / ffmpeg），无法证明视频可播放"
    finally:
        try:
            shutil.rmtree(tmp_dir, ignore_errors=True)
        except Exception:
            pass


def _cover_image_ok(img_bytes: bytes) -> Tuple[bool, str]:
    """封面图检查：
       - 能够被解码器完整加载（image.load() 不抛异常，可"渲染"）；
       - 尺寸有效：宽 ≥ 64 像素、高 ≥ 64 像素（排除占位缩略图/1x1 空图）；
       - 非纯黑：从 25 个采样点中最大亮度 ≥ 30；
       - 非空白：并非所有像素颜色一致（不是纯色占位）。
    """
    try:
        from PIL import Image  # type: ignore
    except ImportError:
        return False, "缺少 PIL/Pillow，无法验证封面清晰度"
    from io import BytesIO
    try:
        img = Image.open(BytesIO(img_bytes))
        img.load()  # 完整解码；损坏图会在此抛异常
        img = img.convert("RGB")
    except Exception as e:
        return False, f"封面无法渲染（解码失败）：{e}"

    w, h = img.size
    if w < 64 or h < 64:
        return False, f"封面尺寸过小 {w}x{h}（<64 像素，疑似占位图）"

    # 25 采样点：5x5 网格
    pts = []
    for i in range(1, 6):
        for j in range(1, 6):
            pts.append((w * i // 6, h * j // 6))
    max_lum = 0.0
    colors = set()
    for x, y in pts:
        px = img.getpixel((x, y))
        if isinstance(px, tuple) and len(px) >= 3:
            r_, g_, b_ = int(px[0]), int(px[1]), int(px[2])
        else:
            v = int(px) if isinstance(px, (int, float)) else 0
            r_ = g_ = b_ = v
        lum = 0.299 * r_ + 0.587 * g_ + 0.114 * b_
        if lum > max_lum:
            max_lum = lum
        colors.add((r_ // 8, g_ // 8, b_ // 8))
    if max_lum < 30:
        return False, f"封面抽样最大亮度 {max_lum:.1f} <30（疑似纯黑）"
    if len(colors) < 3:
        return False, f"封面像素颜色几乎一致（{len(colors)} 种，疑似纯色占位）"
    return True, ""


def check_end_videos_playable(pptx_path: str) -> Tuple[bool, str]:
    """第16页"视频可播放 + 静态封面"细则（每一点都必须满足）：
        播放 (Play):
          PL1) slide16 关系与视频节点结构完整 —— PowerPoint 播放模式启用媒体的必要前提；
          PL2) 点击播放或自动播放触发存在（timing 中 <p:video> 或 hlinkClick media）——
               PowerPoint 播放模式下响应"点击/自动播放"的依据；
          PL3) 引用的视频文件在 pptx 内真实存在（不出现"媒体缺失/路径错误"）；
          PL4) 使用真实媒体解码器（cv2 / ffmpeg CLI）能解码每段视频的首帧，
               证明视频容器与编码可被主流解码器打开，不会在 PowerPoint 中出现"损坏/无法播放"。
          —— PL1~PL4 组合即可覆盖 rubric 中"播放模式下可正常播放、不出现媒体缺失/无法播放/
             路径错误"的判定条件，无需再启用 PowerPoint COM 自动化，避免 Windows +
             PowerPoint 环境依赖以及 SlideShow 窗口副作用。
        静止封面 (Cover):
          CV1) 每段视频均引用了封面 blip；
          CV2/CV3/CV4) 封面可被真实图像解码器完整加载并渲染（Pillow 打开+load 无异常）；
                        尺寸不小于 64x64；不为纯黑；不为空白占位（颜色种类 ≥3）。
    """
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            namelist = set(zf.namelist())

            # 读取 slide16 内容与关系
            if "ppt/slides/slide16.xml" not in namelist:
                return False, "ppt/slides/slide16.xml 不存在"
            slide_xml = zf.read("ppt/slides/slide16.xml").decode("utf-8", errors="ignore")
            rels_path = "ppt/slides/_rels/slide16.xml.rels"
            if rels_path not in namelist:
                return False, "slide16.xml.rels 不存在"
            rels_xml = zf.read(rels_path).decode("utf-8", errors="ignore")

            import re

            # 解析 rels：id -> (type, target)
            rel_map = {}
            for m in re.finditer(
                r'<Relationship\s+[^/]*Id="([^"]+)"[^/]*Type="([^"]+)"[^/]*Target="([^"]+)"',
                rels_xml,
            ):
                rid, rtype, target = m.group(1), m.group(2), m.group(3)
                rel_map[rid] = (rtype, target)

            # PL1) 两段视频对应的 video 节点数量 ≥2
            video_pics = re.findall(r"<p:pic\b[^>]*>.*?</p:pic>", slide_xml, re.DOTALL)
            video_pic_blocks = [b for b in video_pics if "a:videoFile" in b]
            if len(video_pic_blocks) < 2:
                return False, f"slide16 中视频对象 p:pic 数量 {len(video_pic_blocks)} <2"

            # PL2) 点击/自动播放触发
            timing_video_ok = "<p:video>" in slide_xml
            click_play_ok = all("ppaction://media" in b for b in video_pic_blocks)
            if not (timing_video_ok or click_play_ok):
                return False, "未检测到点击播放或自动播放（无 timing video / hlinkClick media）"

            # PL3) 媒体文件存在 & 引用完整
            video_files_seen = set()
            blip_rids = []
            for blk in video_pic_blocks:
                for m in re.finditer(r'r:(?:link|embed)="([^"]+)"', blk):
                    rid = m.group(1)
                    if rid in rel_map:
                        _, target = rel_map[rid]
                        if target.lower().endswith(
                            (".mp4", ".mov", ".avi", ".wmv", ".m4v", ".mkv")
                        ):
                            tgt_norm = (
                                "ppt/" + target.replace("../", "")
                                if target.startswith("..")
                                else target
                            )
                            tgt_norm = tgt_norm.replace("ppt/ppt/", "ppt/")
                            video_files_seen.add(tgt_norm)
                for m in re.finditer(r'<a:blip\s+r:embed="([^"]+)"', blk):
                    blip_rids.append(m.group(1))

            if len(video_files_seen) < 2:
                return False, f"slide16 引用的视频文件数量 {len(video_files_seen)} <2"

            # PL4) 用真实解码器解码每段视频的首帧
            for vf in video_files_seen:
                if vf not in namelist:
                    return False, f"视频文件缺失：{vf}（媒体缺失/路径错误）"
                video_bytes = zf.read(vf)
                ext = os.path.splitext(vf)[1].lower() or ".mp4"
                ok, msg = _decode_first_frame_playable(video_bytes, ext)
                if not ok:
                    return False, f"视频 {vf} 无法解码播放：{msg}"

            # CV1) 每段视频均有封面 blip
            if len(blip_rids) < 2:
                return False, f"封面 blip 数量 {len(blip_rids)} <2，部分视频缺少封面"

            # CV2/CV3/CV4) 封面图像清晰性 —— 用真实图像解码器完整加载并采样
            for rid in blip_rids[:2]:
                if rid not in rel_map:
                    return False, f"封面关系 {rid} 不存在"
                _, target = rel_map[rid]
                tgt = (
                    "ppt/" + target.replace("../", "")
                    if target.startswith("..")
                    else target
                )
                tgt = tgt.replace("ppt/ppt/", "ppt/")
                if tgt not in namelist:
                    return False, f"封面图像缺失：{tgt}"
                img_bytes = zf.read(tgt)
                ok, msg = _cover_image_ok(img_bytes)
                if not ok:
                    return False, f"封面 {tgt} 不合格：{msg}"

    except Exception as e:
        return False, f"读取媒体信息失败：{e}"
    return True, ""


# ---------- 统一入口 ----------
def _locate_pptx(dir_path: str) -> Optional[str]:
    """在给定目录内定位被评估的 .pptx 文档；找不到时返回 None。

    规则：仅识别 .pptx（不再兼容 .ppt 老格式）；忽略 ~$ 打头的临时文件、
    以 . 打头的隐藏文件；同目录多个 .pptx 时按文件名排序取第一个。
    """
    if not os.path.isdir(dir_path):
        return None
    candidates = []
    for name in os.listdir(dir_path):
        if name.startswith("~$") or name.startswith("."):
            continue
        if name.lower().endswith(".pptx"):
            candidates.append(os.path.join(dir_path, name))
    if not candidates:
        return None
    candidates.sort()
    return candidates[0]


def evaluate(dir_path: str) -> dict:
    """脚本对外统一入口：接收"脚本所在目录的路径"，脚本自行在该目录定位并评估文档。

    返回结构（参见"脚本接口差异与统一建议.md" §2.2）：
        {
            "id": "062",
            "file_name": "xxx.pptx",
            "status": "ok" | "error",
            "error": None | str,
            "dim1_pass": bool,
            "dim1_reason": str,
            "dim2_items": [ {rule, max_delta, delta, hit, detail}, ... ],
            "total_score": int,
            "max_score": int,
        }
    """
    result: dict = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": 0,
    }

    try:
        pptx_path = _locate_pptx(dir_path)
        if pptx_path is None:
            result["status"] = "error"
            result["error"] = f"目录 {dir_path} 内未找到 .pptx 文件"
            return result
        result["file_name"] = os.path.basename(pptx_path)

        # 维度一
        dim1_ok, dim1_fails = check_dim1(pptx_path)
        result["dim1_pass"] = bool(dim1_ok)
        result["dim1_reason"] = "" if dim1_ok else "；".join(dim1_fails)

        # 维度二：满分（max_score）固定 = 所有加分项之和（扣分项不计入满分）
        log = Dim2Log()
        if dim1_ok:
            evaluate_dim2(pptx_path, log)
            result["dim2_items"] = log.items
            total = log.total()
            result["total_score"] = max(0, total)
        else:
            # 维度一未通过 —— 不检查维度二，总分 0，dim2_items 保持空
            result["dim2_items"] = []
            result["total_score"] = 0

        # max_score：加分项 max_delta 之和（正值项），保持稳定即便本次 dim2_items 为空
        result["max_score"] = _dim2_max_score()
        return result

    except Exception as e:
        result["status"] = "error"
        result["error"] = f"{type(e).__name__}: {e}\n{traceback.format_exc()}"
        return result


def _dim2_max_score() -> int:
    """维度二加分项的满分总和（用于结果 max_score 字段的稳定输出）。

    与 evaluate_dim2 中的加分项 max_delta 保持一致：
        +5  第2-15页顶部标题外侧蓝色形状组合标题栏
        +5  左侧蓝色主形状
        +5  右侧长条外框 + 斜切过渡 + 线宽
        +5  文本元素（页码/中文/英文）
        +5  装饰点阵 + 短斜线
        +1  第16页"谢谢观看"文本
        +5  第16页视频对象
        +5  视频可播放 + 静态封面
    合计 36 分（扣分项不计入满分）
    """
    return 5 + 5 + 5 + 5 + 5 + 1 + 5 + 5


if __name__ == "__main__":
    # 本地调试入口：evaluate(脚本所在目录) 并打印 JSON
    _here = os.path.dirname(os.path.abspath(__file__))
    _target_dir = sys.argv[1] if len(sys.argv) > 1 else _here
    _r = evaluate(_target_dir)
    print(json.dumps(_r, ensure_ascii=False, indent=2))
