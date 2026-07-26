# -*- coding: utf-8 -*-
"""
对目录内的 PPTX 文件（技术路线图_流程图替换版_调整美化版.pptx）进行自动评估。

严格按打分细则实现 —— 细则中的每一点都对应一段代码；
细则未提的约束一律不加。

评估逻辑：
  维度1 = 交付文件为 .pptx 且能够正常打开（不满足 ⇒ 总分 0，且不再看维度2）
  维度2 = 完成度（6 个 +1、2 个 +3、1 个 +1；扣分项若干）

对外接口（详见「脚本接口差异与统一建议.md §2」）：
    evaluate(dir_path: str) -> dict
        - 入参：脚本所在目录的路径
        - 脚本自己在该目录内定位 .pptx 文件并打开
        - 返回统一的结构化字典（含 dim1_pass / dim2_items / total_score 等）
        - 不 print 主结果、不改 sys.stdout、不 sys.exit、不硬编码路径
"""
from __future__ import annotations

import json
import os
import re
import sys
import zipfile
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.oxml.ns import qn

SCRIPT_ID = "070"
EMU_PER_CM = 360000.0


# =====================================================================
# 基础工具
# =====================================================================

def emu_to_cm(v: Optional[int]) -> Optional[float]:
    return None if v is None else v / EMU_PER_CM


def rect_cm(sh) -> Tuple[float, float, float, float]:
    return (
        emu_to_cm(sh.left) or 0.0,
        emu_to_cm(sh.top) or 0.0,
        emu_to_cm(sh.width) or 0.0,
        emu_to_cm(sh.height) or 0.0,
    )


def rect_contains(outer, inner, tol=0.0) -> bool:
    ol, ot, ow, oh = outer
    il, it_, iw, ih = inner
    return (
        il + tol >= ol
        and it_ + tol >= ot
        and il + iw - tol <= ol + ow
        and it_ + ih - tol <= ot + oh
    )


def rect_overlap(a, b) -> float:
    al, at, aw, ah = a
    bl, bt, bw, bh = b
    ow = max(0.0, min(al + aw, bl + bw) - max(al, bl))
    oh = max(0.0, min(at + ah, bt + bh) - max(at, bt))
    return ow * oh


def shape_text(sh) -> str:
    if not sh.has_text_frame:
        return ""
    return "\n".join(p.text for p in sh.text_frame.paragraphs)


# =====================================================================
# 数据模型
# =====================================================================

@dataclass
class Hit:
    code: str
    points: int
    note: str


@dataclass
class Result:
    dim1_pass: bool = True
    dim1_failures: List[str] = field(default_factory=list)
    hits: List[Hit] = field(default_factory=list)

    @property
    def score(self) -> int:
        if not self.dim1_pass:
            return 0
        return sum(h.points for h in self.hits)


# =====================================================================
# 维度 1：交付文件为 .pptx 且能够正常打开
# =====================================================================

def check_dim1(path: str) -> Tuple[bool, List[str]]:
    fails: List[str] = []

    # 交付文件为.pptx格式，能够正常打开
    if not path.lower().endswith(".pptx"):
        fails.append("文件格式不是 .pptx")
    try:
        with zipfile.ZipFile(path) as zf:
            if not any(n.startswith("ppt/slides/slide") for n in zf.namelist()):
                fails.append("PPTX 内部缺少 slide 部件，无法正常打开")
    except zipfile.BadZipFile:
        fails.append("文件不是有效 OOXML，无法打开")

    return (len(fails) == 0, fails)


# =====================================================================
# 维度 2 辅助：定位 6 个研究板块底框
# 细则把 6 个框按 "下方第一/二/三行" + "左/右列" 给出：
#   01-绿(L,row1)  02-蓝(L,row2)  03-紫(L,row3)
#   04-黄(R,row1)  05-绿(R,row2)  06-粉(R,row3)
# 同时要"位于 '研究方法与技术路线' 内容下方"，所以以该标题文本 top 作为起点。
# =====================================================================

# 颜色族识别：从 SVG 渐变起始 hex 推断
COLOR_FAMILY_RULES = [
    # 顺序敏感：从特征最明显的开始判定
    # yellow: 红绿高、蓝低         （如 FFF1D8）
    ("yellow", lambda r, g, b: r >= 250 and g >= 230 and b < 230),
    # pink:   红高、绿蓝偏低、红>蓝 （如 FFE6EC）
    ("pink",   lambda r, g, b: r >= 250 and g < 240 and b < 245 and r > g and r > b),
    # purple: 红绿差小、蓝最高     （如 F1E9FF）
    ("purple", lambda r, g, b: b >= 250 and r >= 220 and g < r and b > g),
    # blue:   蓝最高、红明显低于蓝 （如 EAF0FF -> 234,240,255; r<b 且 g<b）
    ("blue",   lambda r, g, b: b >= 250 and r < b - 10 and g <= b),
    # green:  绿最高，红蓝接近     （如 E2F8F6, E8F7EA）
    ("green",  lambda r, g, b: g >= 240 and abs(r - b) <= 20 and g >= r),
]


def svg_color_family(svg_bytes: bytes) -> Optional[str]:
    m = re.search(rb'stop-color="#([0-9A-Fa-f]{6})"', svg_bytes)
    if not m:
        return None
    hex6 = m.group(1).decode().upper()
    r, g, b = int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16)
    for fam, pred in COLOR_FAMILY_RULES:
        if pred(r, g, b):
            return fam
    return None


def find_method_title_top(slide) -> Optional[float]:
    """找到 '研究方法与技术路线' 文本框，返回其 top (cm)。"""
    for sh in slide.shapes:
        if sh.has_text_frame and "研究方法与技术路线" in shape_text(sh):
            return rect_cm(sh)[1]
    return None


def collect_section_frames(slide, method_top: float) -> List[dict]:
    """收集 6 个研究板块底框（位于 '研究方法与技术路线' 下方的彩色矩形图片）。"""
    frames = []
    for sh in slide.shapes:
        if sh.shape_type != 13:
            continue
        try:
            blob = sh.image.blob
        except Exception:
            continue
        if not blob.lstrip().startswith(b"<svg"):
            continue
        l, t, ww, hh = rect_cm(sh)
        if t <= method_top:
            continue
        # 6 个底框尺寸基本一致；宽 6~9 cm，高 2~3.5 cm
        if 5.5 < ww < 9.5 and 1.8 < hh < 3.5:
            frames.append({
                "shape": sh,
                "rect": (l, t, ww, hh),
                "color": svg_color_family(blob),
            })
    # 行：按 top 聚类成 3 行；同行按 left 排序
    frames.sort(key=lambda f: f["rect"][1])
    rows: List[List[dict]] = []
    for f in frames:
        if not rows or f["rect"][1] - rows[-1][0]["rect"][1] > 1.0:
            rows.append([f])
        else:
            rows[-1].append(f)
    for row in rows:
        row.sort(key=lambda f: f["rect"][0])
    return rows  # type: ignore


# 细则指定：(板块编号key, 标题原文完整串, 行号(0-based), 期望颜色)
# 列(L/R)不是细则约束 —— 在该行内按颜色定位框。
SECTIONS = [
    ("研究内容一", "研究内容一：关键病理事件界定",   0, "green"),
    ("研究内容二", "研究内容二：色氨酸代谢异常定位", 1, "blue"),
    ("研究内容三", "研究内容三：ATF4核内调控解析",   2, "purple"),
    ("研究内容四", "研究内容四：表观遗传持续化机制", 0, "yellow"),
    ("研究内容五", "研究内容五：轴向因果验证",       1, "green"),
    ("预期结果", "预期结果与科学意义",             2, "pink"),
]


# 细则「4 项正文」的原文（按板块 key 索引）。有原文的板块用逐项匹配，没有的板块
# 回退到条目计数。原文来源：任务书 / 细则里给出的板块内容。
EXPECTED_BODY_ITEMS: Dict[str, List[str]] = {
    "研究内容二": [
        "卵泡液与颗粒细胞代谢组检测",
        "筛选差异代谢物并锁定关键节点",
        "重点观察：IA及相关吲哚类代谢物",
        "建立代谢物—黄体化表型关联",
    ],
    "研究内容三": [
        "ATF4表达、磷酸化与核定位检测",
        "核/胞分离与免疫荧光验证",
        "ChIP-qPCR / CUT&Tag定位靶基因",
        "解析“代谢异常—核内转录”桥梁",
    ],
    "研究内容四": [
        "H3K4me3与H3K27ac开放状态检测",
        "关键位点：STAR / COX2 / PTX3",
        "刺激撤除后恢复观察",
        "判断程序异常是否具备持续化特征",
    ],
    "研究内容五": [
        "IA梯度干预与时间窗干预",
        "ATF4沉默/过表达与救援实验",
        "表观修饰干预与逆转观察",
        "形成“IA—ATF4—表观遗传”证据链",
    ],
    "预期结果": [
        "揭示PCOS持续性排卵障碍的新病理机制",
        "建立“代谢异常—核内调控—程序启动失败”框架",
        "提出可分层评估的促排反应预测标志",
        "为靶向代谢—表观干预提供实验依据",
    ],
}


# 部分板块细则里额外约束了"编号 + 页面左右位置"（如"右侧04"）。
# 结构：板块 key -> (期望编号文本, 期望页面列 "L"/"R")
EXPECTED_SECTION_POS: Dict[str, Tuple[str, str]] = {
    "研究内容四": ("04", "R"),
    "研究内容五": ("05", "R"),
}


def _norm_body_text(s: str) -> str:
    """归一化正文字符串：去空白 + 统一中英文标点 + 去除项目符号/编号前缀。

    用于逐项子串匹配，避免全/半角、空格、序号（"1." / "1、"）导致误判。
    """
    if not s:
        return ""
    # 中英文标点统一
    trans = str.maketrans({
        "：": ":", "，": ",", "。": ".", "；": ";",
        "（": "(", "）": ")", "—": "-", "–": "-",
        "、": ",",
        "“": '"', "”": '"', "‘": "'", "’": "'",
    })
    x = s.translate(trans)
    # 去掉行首"1." / "1)" / "1、" / "①" 等序号前缀
    x = re.sub(r"^\s*(?:\d+\s*[.)、:]|[①-⑳])", "", x)
    # 去掉常见项目符号
    x = re.sub(r"[·•●◆■▶►\s]", "", x)
    return x


def _body_items_match(body_sh, expected: List[str]) -> bool:
    """判定正文中"逐项"都出现（归一化后子串匹配）。

    - 拆条目：与 _body_item_count 一致（按段落 + 项目符号/分隔符）；
    - 每条 expected 归一化后必须能作为子串出现在"任一实际条目的归一化文本"中；
    - 每条 expected 只允许被消费一次，避免同一条目重复计数；
    - 全部命中才通过。
    """
    if body_sh is None or not expected:
        return False
    substantive = re.compile(r"[一-鿿A-Za-z0-9]")
    splitter = re.compile(r"[\n；;•·●◆■▶►]")
    actual_norm: List[str] = []
    for p in body_sh.text_frame.paragraphs:
        raw = p.text or ""
        for seg in splitter.split(raw):
            s = seg.strip(" \t\r　·•●◆■▶►、,，.。;；:：-—")
            if len(s) >= 2 and substantive.search(s):
                actual_norm.append(_norm_body_text(s))
    used = [False] * len(actual_norm)
    for exp in expected:
        exp_n = _norm_body_text(exp)
        if not exp_n:
            return False
        matched = False
        for i, a in enumerate(actual_norm):
            if used[i]:
                continue
            if exp_n in a or a in exp_n:
                used[i] = True
                matched = True
                break
        if not matched:
            return False
    return True


def _body_content_fits(body_sh) -> bool:
    """判定正文文本框内的内容是否完整显示（未被裁切、未被自动收缩）。

    仅靠"文本框矩形 ⊆ 底框矩形"不够——PPT 允许文本溢出文本框自身，或开启
    "缩排以适应"(normAutofit) 让文字自动缩字。这里补齐三道判据：
      1) 文本框 wrap=off ⇒ 直接判失败（默认开启自动换行；关闭意味着允许水平溢出）；
      2) <a:normAutofit fontScale="…"> 存在且 fontScale < 100000 ⇒ PPT 触发了
         自动缩字，说明按名义字号已经放不下，属于隐性裁切/收缩；
      3) 估算文本堆叠高度（行数 × 字号 × 行距 × 1.2 pt→cm 换算）不得超过
         文本框可用高度（扣除上下 0.13cm 默认内边距 * 2）。
    """
    if body_sh is None or not body_sh.has_text_frame:
        return True
    tf = body_sh.text_frame

    # (1) 自动换行关闭 ⇒ 允许水平溢出
    if getattr(tf, "word_wrap", None) is False:
        return False

    # (2) normAutofit 缩字系数：<a:bodyPr><a:normAutofit fontScale="90000"/></a:bodyPr>
    try:
        bodyPr = tf._txBody.find(qn("a:bodyPr"))
        if bodyPr is not None:
            na = bodyPr.find(qn("a:normAutofit"))
            if na is not None:
                fs = na.get("fontScale")
                if fs is not None and int(fs) < 100000:
                    return False
                lnSpcReduction = na.get("lnSpcReduction")
                if lnSpcReduction is not None and int(lnSpcReduction) > 0:
                    return False
    except Exception:
        pass

    # (3) 估算堆叠高度
    _, _, _, bh_cm = rect_cm(body_sh)
    if bh_cm <= 0:
        return True
    # 上下 inset：默认约 0.13cm (0.05")；若显式设置则取实际值
    inset_top_cm = 0.13
    inset_bot_cm = 0.13
    try:
        if bodyPr is not None:
            tIns = bodyPr.get("tIns")
            bIns = bodyPr.get("bIns")
            if tIns is not None:
                inset_top_cm = int(tIns) / EMU_PER_CM
            if bIns is not None:
                inset_bot_cm = int(bIns) / EMU_PER_CM
    except Exception:
        pass
    available_cm = bh_cm - inset_top_cm - inset_bot_cm
    if available_cm <= 0:
        return False

    PT_TO_CM = 2.54 / 72.0
    stack_cm = 0.0
    for p in tf.paragraphs:
        # 段前 / 段后
        try:
            sb_pt = p.space_before.pt if p.space_before is not None else 0.0
        except Exception:
            sb_pt = 0.0
        try:
            sa_pt = p.space_after.pt if p.space_after is not None else 0.0
        except Exception:
            sa_pt = 0.0
        # 行高：取该段内 run 的最大 font.size；缺省按 8pt 估计（细则正文范围）
        run_sizes = [r.font.size.pt for r in p.runs if r.font.size is not None]
        size_pt = max(run_sizes) if run_sizes else 8.0
        # 行距倍数：float 直接用；None/其它当 1.0；绝对 pt 值超出细则范畴按 1.0 估
        ls = p.line_spacing
        if isinstance(ls, float) and ls > 0:
            ls_mul = ls
        else:
            ls_mul = 1.0
        # 单行高（pt）≈ 字号 × 行距 × 1.2；每段至少 1 行（空行也占位）
        line_pt = size_pt * ls_mul * 1.2
        stack_cm += (sb_pt + line_pt + sa_pt) * PT_TO_CM

    # 3% 浮点/字体度量容差
    return stack_cm <= available_cm * 1.03


def locate_frame(rows: List[List[dict]], row_idx: int, color: str) -> Optional[dict]:
    """在指定行内按颜色定位底框（细则只给了行号+颜色，没给左/右）。"""
    if row_idx >= len(rows) or not rows[row_idx]:
        return None
    for f in rows[row_idx]:
        if f["color"] == color:
            return f
    return None


def shapes_inside_frame(slide, frame_rect: Tuple[float, float, float, float]) -> List:
    """返回完全位于 frame_rect 内的可编辑文本 shape（容差 0.20cm）。

    仅收集完整位于目标框内的候选，避免把未完整放入框内的文本参与后续
    "标题 / 正文 / 编号" 分类。
    """
    inside = []
    for sh in slide.shapes:
        if not sh.has_text_frame or not shape_text(sh).strip():
            continue
        r = rect_cm(sh)
        if r[2] <= 0 or r[3] <= 0:
            continue
        if rect_contains(frame_rect, r, tol=0.20):
            inside.append(sh)
    return inside


def _body_item_count(body_sh) -> int:
    """统计正文中"实际条目"数量。

    条目 = text_frame 里"非空且有实质内容"的段落：
      - 去除首尾空白后长度 >= 2；
      - 至少包含一个汉字或字母数字字符（避免纯标点/纯符号误当条目）；
      - 若一段中出现明显的分隔符（换行以外的项目标识、如全角分号、• ），
        则按其拆分再逐一校验，兼容"多项写在同一段"的写法。
        注意：不用中文顿号"、"作为条目分隔符——它更常出现在条目内部
        （如"表达、磷酸化"）。
    """
    if body_sh is None:
        return 0
    substantive = re.compile(r"[一-鿿A-Za-z0-9]")
    splitter = re.compile(r"[\n；;•·●◆■▶►]")
    count = 0
    for p in body_sh.text_frame.paragraphs:
        raw = p.text or ""
        for seg in splitter.split(raw):
            s = seg.strip(" \t\r　·•●◆■▶►、,，.。;；:：-—")
            if len(s) >= 2 and substantive.search(s):
                count += 1
    return count


def _norm_title(s: str) -> str:
    """规整化：去掉空白，统一中英文冒号。"""
    return re.sub(r"\s+", "", s or "").replace(":", "：")


# 细则给出的 6 个板块标题原文（用于识别哪个 shape 是"标题"）
SECTION_TITLE_TEXTS = {
    "研究内容一：关键病理事件界定",
    "研究内容二：色氨酸代谢异常定位",
    "研究内容三：ATF4核内调控解析",
    "研究内容四：表观遗传持续化机制",
    "研究内容五：轴向因果验证",
    "预期结果与科学意义",
}
_NORM_TITLES = {_norm_title(t) for t in SECTION_TITLE_TEXTS}


def classify_section_shapes(shapes) -> Dict[str, Optional[object]]:
    """从框内 shape 中分离出 编号 / 标题 / 正文。
    标题：与细则给出的 6 条标题原文之一字面匹配（忽略空白与冒号全/半角）。
    编号：0[1-6]。
    正文：其它含文本的 shape，取最长者。
    """
    num = title = body = None
    for sh in shapes:
        txt = shape_text(sh).strip()
        if re.fullmatch(r"0[1-6]", txt):
            num = sh
            continue
        norm = _norm_title(txt)
        if any(t in norm for t in _NORM_TITLES):
            title = sh
            continue
        if body is None or len(shape_text(sh)) > len(shape_text(body)):
            body = sh
    return {"num": num, "title": title, "body": body}


# =====================================================================
# 维度 2 主体
# =====================================================================

def evaluate_dim2(slide, prs) -> List[Hit]:
    hits: List[Hit] = []

    method_top = find_method_title_top(slide)
    if method_top is None:
        # 没找到 "研究方法与技术路线" 标题：6 个 +1 全部无法判定，直接返回
        return hits

    rows = collect_section_frames(slide, method_top)

    # ---------- 6 条 +1：每个研究板块标题+4项正文完整放入对应颜色框，不超出 ----------
    section_info = []   # 后面 +3/+1/扣分 都要用

    for code_key, full_title, row_idx, want_color in SECTIONS:
        frame = locate_frame(rows, row_idx, want_color)
        rec = {"key": code_key, "full_title": full_title, "row": row_idx,
               "want_color": want_color, "frame": frame,
               "title": None, "body": None, "num": None,
               "inside": False, "color_ok": False, "body_4": False,
               "pos_ok": True,
               "col": None}
        if frame is None:
            section_info.append(rec)
            continue
        # 页面左右列判定：以整页宽度的中线为准（w_cm 为 None 时回退到 5.0 阈值）
        w_cm = emu_to_cm(prs.slide_width) or 20.0
        frame_center_x = frame["rect"][0] + frame["rect"][2] / 2.0
        rec["col"] = "L" if frame_center_x < w_cm / 2.0 else "R"
        parts = classify_section_shapes(shapes_inside_frame(slide, frame["rect"]))
        title_sh, body_sh, num_sh = parts["title"], parts["body"], parts["num"]
        rec["title"], rec["body"], rec["num"] = title_sh, body_sh, num_sh

        # 标题：与细则给出的标题原文逐字匹配（忽略空白/冒号全半角差异）
        title_match = (title_sh is not None
                       and _norm_title(full_title) in _norm_title(shape_text(title_sh)))
        # "4 项正文"：有原文的板块按逐项匹配；未配置的板块回退到"条目数==4"计数。
        expected_items = EXPECTED_BODY_ITEMS.get(code_key)
        if expected_items is not None:
            rec["body_4"] = _body_items_match(body_sh, expected_items)
        else:
            rec["body_4"] = (body_sh is not None and _body_item_count(body_sh) == 4)

        # 细则"不超出框体"：标题与正文都必须被底框包住（容差 0.20cm，浮点抖动用）
        inside_ok = True
        if title_sh is not None:
            inside_ok &= rect_contains(frame["rect"], rect_cm(title_sh), tol=0.20)
        if body_sh is not None:
            inside_ok &= rect_contains(frame["rect"], rect_cm(body_sh), tol=0.20)
        rec["inside"] = inside_ok
        # locate_frame 已经按颜色筛选，这里冗余确认
        rec["color_ok"] = (frame["color"] == want_color)

        # 细则里额外约束了"编号 + 左右位置"的板块（如"右侧04"），做硬校验
        pos_req = EXPECTED_SECTION_POS.get(code_key)
        if pos_req is not None:
            want_num, want_col = pos_req
            num_text = shape_text(num_sh).strip() if num_sh is not None else ""
            rec["pos_ok"] = (num_text == want_num and rec["col"] == want_col)

        section_info.append(rec)

        if title_match and rec["body_4"] and inside_ok and rec["color_ok"] and rec["pos_ok"]:
            rules = {
                "研究内容一": "“研究内容一：关键病理事件界定”板块：标题和4项正文完整放入”研究方法与技术路线”内容下方第一行的的绿色框内，不超出框体。",
                "研究内容二": "“研究内容二：色氨酸代谢异常定位”板块：标题和4项正文完整放入”研究方法与技术路线”内容下方第二行蓝色框内，不超出框体。",
                "研究内容三": "“研究内容三：ATF4核内调控解析”板块：标题和4项正文完整放入”研究方法与技术路线”内容下方第三行的紫色框内,不超出框体。",
                "研究内容四": "“研究内容四：表观遗传持续化机制”板块：标题和4项正文完整放入”研究方法与技术路线”内容下方的第一行的黄色框内，不超出框体。",
                "研究内容五": "“研究内容五：轴向因果验证”板块：标题和4项正文完整放入”研究方法与技术路线”内容下方第二行的绿色框内，不超出框体。",
                "预期结果":   "“预期结果与科学意义”板块：标题和4项正文完整放入”研究方法与技术路线”内容下方第三行的粉色框内，不超出框体。",
            }
            hits.append(Hit(f"+1#{code_key}", 1, rules[code_key]))

    # ---------- +3：六个研究板块标题统一使用约10.5至11.5磅、加粗、左对齐，位于框内顶部；
    #              正文统一使用约7.2至8.2磅、常规字重，字号能够保证全部内容完整显示。 ----------
    # 严格按 rubric 范围，仅给 ±0.05pt 极小浮点容差（PPT 存储 1/100 pt，实际数值抖动很小）。
    TITLE_MIN, TITLE_MAX = 10.5 - 0.05, 11.5 + 0.05
    BODY_MIN,  BODY_MAX  = 7.2 - 0.05,  8.2 + 0.05

    title_size_ok = True   # 约10.5–11.5pt
    title_bold_ok = True   # 加粗
    title_left_ok = True   # 左对齐
    title_top_ok  = True   # 位于框内顶部
    body_size_ok  = True   # 约7.2–8.2pt
    body_normal_ok = True  # 常规字重
    body_fit_ok    = True  # 字号能够保证全部内容完整显示（不溢出框体）
    have_titles = have_bodies = 0

    for rec in section_info:
        t, b, f = rec["title"], rec["body"], rec["frame"]
        if f is None:
            continue
        if t is not None:
            have_titles += 1
            sizes = [r.font.size.pt for p in t.text_frame.paragraphs for r in p.runs if r.font.size]
            if not sizes or not all(TITLE_MIN <= s <= TITLE_MAX for s in sizes):
                title_size_ok = False
            bolds = [bool(r.font.bold) for p in t.text_frame.paragraphs for r in p.runs]
            if not bolds or not all(bolds):
                title_bold_ok = False
            # 左对齐：PP_ALIGN.LEFT 或 None（默认左对齐）均视为合规
            aligns = [p.alignment for p in t.text_frame.paragraphs]
            if any(a is not None and a != PP_ALIGN.LEFT for a in aligns):
                title_left_ok = False
            # 位于框内顶部：标题 top 距框 top ≤ 框高的 30%
            f_top = f["rect"][1]
            f_h   = f["rect"][3]
            if (rect_cm(t)[1] - f_top) > max(0.5, f_h * 0.30):
                title_top_ok = False
        if b is not None:
            have_bodies += 1
            sizes = [r.font.size.pt for p in b.text_frame.paragraphs for r in p.runs if r.font.size]
            if not sizes or not all(BODY_MIN <= s <= BODY_MAX for s in sizes):
                body_size_ok = False
            bolds = [bool(r.font.bold) for p in b.text_frame.paragraphs for r in p.runs]
            if any(bolds):
                body_normal_ok = False
            # 字号能保证全部内容完整显示：多重判据
            #   (a) 正文矩形被底框完整包住（外部溢出）；
            #   (b) 文本框未设置"缩排以适应"(normAutofit) 收缩系数（<1 表示 PPT 自动缩字，
            #       说明按当前字号已放不下，属于隐性裁切/收缩）；
            #   (c) 关闭 word_wrap 时视为"手动允许溢出"，直接判失败；
            #   (d) 估算文本堆叠高度（行数 × 字号 × 行距 × 1.2 换算 cm）不得超过
            #       正文框可用高度（扣除上下 0.2cm 边距），超过判失败。
            if not rect_contains(f["rect"], rect_cm(b), tol=0.20):
                body_fit_ok = False
            if not _body_content_fits(b):
                body_fit_ok = False

    if (have_titles == 6 and have_bodies == 6
            and title_size_ok and title_bold_ok and title_left_ok and title_top_ok
            and body_size_ok and body_normal_ok and body_fit_ok):
        hits.append(Hit(
            "+3#font", 3,
            "六个研究板块标题统一使用约10.5至11.5磅、加粗、左对齐，位于框内顶部。正文统一使用约7.2至8.2磅、常规字重，字号能够保证全部内容完整显示。"
        ))

    # ---------- +3：六个研究板块正文行距：统一设置为单倍或约0.9至1.0倍行距，
    #              段前、段后间距为0至2磅。文字距框体左右边线约0.3至0.5厘米，
    #              距上下边线不少于0.2厘米。 ----------
    # "约 0.3–0.5cm" 给 ±0.05cm 容差（细则用了"约"字）。
    LR_MIN, LR_MAX = 0.3 - 0.05, 0.5 + 0.05

    line_ok = before_ok = after_ok = margin_ok = True
    have_bodies_with_frame = 0
    for rec in section_info:
        b, f = rec["body"], rec["frame"]
        if b is None or f is None:
            continue
        have_bodies_with_frame += 1
        for p in b.text_frame.paragraphs:
            ls = p.line_spacing
            # 单倍 或 约 0.9–1.0 倍；None 即 PPT 默认（单倍）
            if ls is None:
                pass
            elif isinstance(ls, float):
                if not (0.9 <= ls <= 1.0):
                    line_ok = False
            else:
                # 绝对值（pt）不属于细则"倍数"体系
                line_ok = False
            # 段前 0–2 磅；None 视为 0
            sb = p.space_before.pt if p.space_before is not None else 0.0
            if not (0.0 <= sb <= 2.0):
                before_ok = False
            # 段后 0–2 磅；None 视为 0
            sa = p.space_after.pt if p.space_after is not None else 0.0
            if not (0.0 <= sa <= 2.0):
                after_ok = False
        fl, ft, fw, fh = f["rect"]
        bl, bt, bw, bh = rect_cm(b)
        left_m  = bl - fl
        right_m = (fl + fw) - (bl + bw)
        top_m   = bt - ft
        bot_m   = (ft + fh) - (bt + bh)
        # 左右：约 0.3–0.5cm
        if not (LR_MIN <= left_m <= LR_MAX and LR_MIN <= right_m <= LR_MAX):
            margin_ok = False
        # 上下：不少于 0.2cm
        if not (top_m >= 0.2 and bot_m >= 0.2):
            margin_ok = False

    if have_bodies_with_frame == 6 and line_ok and before_ok and after_ok and margin_ok:
        hits.append(Hit(
            "+3#spacing", 3,
            "六个研究板块正文行距：统一设置为单倍或约0.9至1.0倍行距，段前、段后间距为0至2磅。文字距框体左右边线约0.3至0.5厘米，距上下边线不少于0.2厘米。"
        ))

    # ---------- +1：左侧01、02、03板块和右侧04、05、06板块标题、正文和编号左边缘对齐，
    #              上下三块保持相同排版规则。 ----------
    # 细则把六个板块明确指名：01/02/03 在左、04/05/06 在右。直接用编号文本归类，
    # 不依赖前面"col"启发式。

    def _aligned(vals: List[float], tol=0.15) -> bool:
        # 同列三块左边缘"对齐"：极差容差 0.15cm
        return len(vals) == 3 and (max(vals) - min(vals)) <= tol

    left_titles, left_bodies, left_nums = [], [], []
    right_titles, right_bodies, right_nums = [], [], []

    for rec in section_info:
        if rec["frame"] is None or rec["num"] is None:
            continue
        num_text = shape_text(rec["num"]).strip()
        if num_text in ("01", "02", "03"):
            target_t, target_b, target_n = left_titles, left_bodies, left_nums
        elif num_text in ("04", "05", "06"):
            target_t, target_b, target_n = right_titles, right_bodies, right_nums
        else:
            continue
        if rec["title"] is not None:
            target_t.append(rect_cm(rec["title"])[0])
        if rec["body"] is not None:
            target_b.append(rect_cm(rec["body"])[0])
        target_n.append(rect_cm(rec["num"])[0])

    # 左侧三块、右侧三块的 标题/正文/编号 各自左边缘必须对齐
    left_ok = _aligned(left_titles) and _aligned(left_bodies) and _aligned(left_nums)
    right_ok = _aligned(right_titles) and _aligned(right_bodies) and _aligned(right_nums)

    # "上下三块保持相同排版规则"：左右两列内部，三块的 (标题-编号, 正文-编号)
    # 偏移在 ±0.15cm 内一致；同时左右两列之间的偏移模式也保持一致（±0.20cm）。
    def _const(vals: List[float], tol=0.15) -> bool:
        return len(vals) == 3 and (max(vals) - min(vals)) <= tol

    rule_same = True
    # 在每一列内：三块的"标题相对编号""正文相对编号"偏移应一致
    for nums, titles, bodies in [
        (left_nums, left_titles, left_bodies),
        (right_nums, right_titles, right_bodies),
    ]:
        if len(nums) == 3 and len(titles) == 3 and len(bodies) == 3:
            t_offsets = [t - n for t, n in zip(titles, nums)]
            b_offsets = [b - n for b, n in zip(bodies, nums)]
            if not (_const(t_offsets) and _const(b_offsets)):
                rule_same = False
        else:
            rule_same = False
    # 左右两列之间排版规则一致：偏移均值差 < 0.20cm
    if rule_same and left_titles and right_titles and left_nums and right_nums:
        l_off = sum(t - n for t, n in zip(left_titles, left_nums)) / 3
        r_off = sum(t - n for t, n in zip(right_titles, right_nums)) / 3
        if abs(l_off - r_off) > 0.20:
            rule_same = False

    if left_ok and right_ok and rule_same:
        hits.append(Hit(
            "+1#align", 1,
            "左侧01、02、03板块和右侧04、05、06板块标题、正文和编号左边缘对齐，上下三块保持相同排版规则"
        ))

    # ============================ 扣分项 ============================

    w_cm = emu_to_cm(prs.slide_width)
    h_cm = emu_to_cm(prs.slide_height)

    # -1：页面尺寸不是 20.00×28.29 厘米
    size_ok = (w_cm and h_cm and abs(w_cm - 20.00) <= 0.05 and abs(h_cm - 28.29) <= 0.05)
    if not size_ok:
        hits.append(Hit("-1#size", -1, "页面尺寸不是20.00×28.29厘米"))

    return hits


# =====================================================================
# 报告 / 统一接口
# =====================================================================

# 维度 2 全部评分项：命中与未命中都要在返回结果中列出（顺序即输出顺序）。
# 结构：(code, max_delta, penalty_delta, rule_text)
#   - 加分项：max_delta = 满分，penalty_delta = 0
#   - 扣分项：max_delta = 触发后的负分，penalty_delta 与 max_delta 一致
DIM2_RULES: List[Tuple[str, int, int, str]] = [
    ("+1#研究内容一", 1, 0,
     "“研究内容一：关键病理事件界定”板块：标题和4项正文完整放入”研究方法与技术路线”内容下方第一行的的绿色框内，不超出框体。"),
    ("+1#研究内容二", 1, 0,
     "“研究内容二：色氨酸代谢异常定位”板块：标题和4项正文完整放入”研究方法与技术路线”内容下方第二行蓝色框内，不超出框体。"),
    ("+1#研究内容三", 1, 0,
     "“研究内容三：ATF4核内调控解析”板块：标题和4项正文完整放入”研究方法与技术路线”内容下方第三行的紫色框内,不超出框体。"),
    ("+1#研究内容四", 1, 0,
     "“研究内容四：表观遗传持续化机制”板块：标题和4项正文完整放入”研究方法与技术路线”内容下方的第一行的黄色框内,不超出框体。"),
    ("+1#研究内容五", 1, 0,
     "“研究内容五：轴向因果验证”板块：标题和4项正文完整放入”研究方法与技术路线”内容下方第二行的绿色框内，不超出框体。"),
    ("+1#预期结果", 1, 0,
     "“预期结果与科学意义”板块：标题和4项正文完整放入”研究方法与技术路线”内容下方第三行的粉色框内，不超出框体。"),
    ("+3#font", 3, 0,
     "六个研究板块标题统一使用约10.5至11.5磅、加粗、左对齐，位于框内顶部。正文统一使用约7.2至8.2磅、常规字重，字号能够保证全部内容完整显示。"),
    ("+3#spacing", 3, 0,
     "六个研究板块正文行距：统一设置为单倍或约0.9至1.0倍行距，段前、段后间距为0至2磅。文字距框体左右边线约0.3至0.5厘米，距上下边线不少于0.2厘米。"),
    ("+1#align", 1, 0,
     "左侧01、02、03板块和右侧04、05、06板块标题、正文和编号左边缘对齐，上下三块保持相同排版规则。"),
    ("-1#size", -1, -1, "页面尺寸不是20.00×28.29厘米。"),
]


def _find_pptx_in_dir(dir_path: str) -> Optional[str]:
    """在给定目录内定位 .pptx 文件；若有多个，取修改时间最新者。"""
    if not os.path.isdir(dir_path):
        return None
    candidates: List[str] = []
    for name in os.listdir(dir_path):
        low = name.lower()
        if low.endswith(".pptx") and not name.startswith("~$"):
            candidates.append(os.path.join(dir_path, name))
    if not candidates:
        return None
    candidates.sort(key=lambda p: os.path.getmtime(p), reverse=True)
    return candidates[0]


def _run_evaluation(path: str) -> Result:
    res = Result()
    prs = Presentation(path)
    ok, fails = check_dim1(path)
    res.dim1_pass = ok
    res.dim1_failures = fails
    if not ok:
        return res
    res.hits = evaluate_dim2(prs.slides[0], prs)
    return res


def _build_dim2_items(res: Result) -> List[Dict[str, object]]:
    hit_map = {h.code: h for h in res.hits}
    items: List[Dict[str, object]] = []
    for code, max_delta, _penalty_delta, rule in DIM2_RULES:
        h = hit_map.get(code)
        hit = h is not None
        if hit:
            delta = h.points
            detail = ""
        else:
            delta = 0
            detail = ""
        items.append({
            "rule": rule,
            "max_delta": max_delta,
            "delta": delta,
            "hit": hit,
            "detail": detail,
        })
    return items


def evaluate(dir_path: str) -> Dict[str, object]:
    """统一入口：接收脚本所在目录路径，脚本自行在目录内定位并打开被评估文档。

    返回结构见「脚本接口差异与统一建议.md §2.2」。
    """
    result: Dict[str, object] = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": sum(max_delta for _, max_delta, _, _ in DIM2_RULES if max_delta > 0),
    }
    try:
        target = _find_pptx_in_dir(dir_path)
        if target is None:
            result["status"] = "error"
            result["error"] = f"目录内未找到 .pptx 文件：{dir_path}"
            return result
        result["file_name"] = os.path.basename(target)

        res = _run_evaluation(target)
        result["dim1_pass"] = res.dim1_pass
        if not res.dim1_pass:
            result["dim1_reason"] = "；".join(res.dim1_failures)
            result["dim2_items"] = []
            result["total_score"] = 0
            return result
        result["dim2_items"] = _build_dim2_items(res)
        result["total_score"] = res.score
        return result
    except Exception as e:
        result["status"] = "error"
        result["error"] = f"{type(e).__name__}: {e}"
        return result


if __name__ == "__main__":
    # 本地调试用：默认取脚本所在目录；也可通过 argv[1] 显式指定
    target_dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(target_dir), ensure_ascii=False, indent=2))
