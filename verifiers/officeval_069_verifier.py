# -*- coding: utf-8 -*-
"""
对 "工商层面_流程图_框体字体调整版.pptx" 的自动评估脚本

评估逻辑：
  维度1（交付文件为 .pptx 格式，能够正常打开）：不通过 -> 0 分，且不再检查维度2。
  维度2（完成度）：累计加分项和扣分项。
    - 加分项必须满足该项中的每一个细则才得分；
    - 扣分项只要满足该项中的任意一个细则即扣分。

对外接口：
    def evaluate(dir_path: str) -> dict
        - dir_path：脚本所在目录的路径；脚本自行在该目录中定位并打开被评估文档。
        - 返回结构化字典，字段见文件末尾示例。
"""

import os
import sys

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- 工具函数 ----------

EMU_PER_CM = 360000.0

def emu_to_cm(v):
    return v / EMU_PER_CM if v is not None else None

def approx(value_cm, target_cm, tol=0.25):
    """容差比较（默认 ±0.25 cm，约 7 emu·万）"""
    if value_cm is None:
        return False
    return abs(value_cm - target_cm) <= tol

def hex_eq(a, b):
    if a is None or b is None:
        return False
    return str(a).upper().lstrip("#") == str(b).upper().lstrip("#")

def color_close(a, b, tol=24):
    """颜色容差比较：a/b 为 6 位 HEX 字符串"""
    if a is None or b is None:
        return False
    try:
        a = str(a).upper().lstrip("#")
        b = str(b).upper().lstrip("#")
        ar, ag, ab = int(a[0:2], 16), int(a[2:4], 16), int(a[4:6], 16)
        br, bg, bb = int(b[0:2], 16), int(b[2:4], 16), int(b[4:6], 16)
        return abs(ar - br) <= tol and abs(ag - bg) <= tol and abs(ab - bb) <= tol
    except Exception:
        return False

def get_shape_fill_hex(shape):
    try:
        f = shape.fill
        if f.type == 1:  # SOLID
            return str(f.fore_color.rgb)
    except Exception:
        pass
    return None

def get_shape_line_hex(shape):
    try:
        return str(shape.line.color.rgb)
    except Exception:
        return None

def get_font_summary(shape):
    """汇总文本框内字体属性。返回 dict：names/sizes/bolds/aligns/vanchor"""
    info = {"names": set(), "sizes": [], "bolds": [], "aligns": set(), "vanchor": None}
    if not shape.has_text_frame:
        return info
    try:
        info["vanchor"] = shape.text_frame.vertical_anchor
    except Exception:
        info["vanchor"] = None
    for para in shape.text_frame.paragraphs:
        if para.alignment is not None:
            info["aligns"].add(para.alignment)
        for run in para.runs:
            if run.text.strip() == "":
                continue
            if run.font.name:
                info["names"].add(run.font.name)
            if run.font.size is not None:
                info["sizes"].append(run.font.size.pt)
            info["bolds"].append(bool(run.font.bold))
    return info

def font_ok(info, want_name, want_pt, want_bold=True, pt_tol=1.5):
    if want_name not in info["names"]:
        return False
    if want_bold and not all(info["bolds"]) if info["bolds"] else (want_bold is False):
        # 全部 run 都要 bold
        if want_bold and not all(info["bolds"]):
            return False
    if not info["sizes"]:
        return False
    avg_pt = sum(info["sizes"]) / len(info["sizes"])
    if abs(avg_pt - want_pt) > pt_tol:
        return False
    return True

def alignment_ok(info):
    """水平 + 垂直居中"""
    h_ok = (len(info["aligns"]) == 0) or (info["aligns"] == {PP_ALIGN.CENTER}) or (PP_ALIGN.CENTER in info["aligns"] and len(info["aligns"]) == 1)
    v_ok = info["vanchor"] == MSO_ANCHOR.MIDDLE
    return h_ok and v_ok

def find_shape_by_text(slide, keywords_all):
    """按关键字（必须全部包含）查找形状"""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text
        if all(k in txt for k in keywords_all):
            return sh
    return None

def is_rounded_rect(shape):
    try:
        return "ROUNDED_RECTANGLE" in str(shape.auto_shape_type)
    except Exception:
        return False


# ---------- 维度1 检查 ----------

def check_dimension_1(pptx_path):
    """返回 (passed: bool, details: list[str])

    细则：交付文件为 .pptx 格式，能够正常打开。
    """
    details = []
    # 1) 格式：文件后缀为 .pptx
    if not pptx_path.lower().endswith(".pptx"):
        return False, ["文件后缀不是 .pptx"]
    # 2) 能够正常打开：python-pptx 解析成功
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return False, [f"无法打开 pptx：{e}"]
    details.append("✓ 文件可打开（.pptx 格式，python-pptx 解析成功）")

    if len(prs.slides) == 0:
        return False, details + ["✗ PPT 为空"]

    return True, details


# ---------- 维度2 评分细则定义 ----------
# 每个加分项是一个统一的"模板"：
#   关键字（必须全部包含的子串） / 期望宽高（cm） / 填充 HEX / 框线 HEX / 字体名 / 字号 pt / 是否加粗 / 是否水平垂直居中

PLUS_RULES = [
    # 标题下方两个时间轴
    dict(score=3,
         label='"2023年 管理负责人调整 职能重新划分"外框：宽约5.08厘米、高约2.11厘米的圆角矩形，使用浅蓝色背景，颜色代码为DDE9FF，深蓝色框线，颜色代码为3768B3，字体为Noto Sans CJK SC、约13.2磅、加粗、水平及垂直居中。',
         keys=["2023年", "管理负责人调整", "职能重新划分"],
         w=5.08, h=2.11, fill="DDE9FF", line="3768B3",
         font="Noto Sans CJK SC", pt=13.2, bold=True, center=True,
         use_strict_rule=True),
    dict(score=3,
         label='"2025年 引入合规与数据 治理机制"外框：宽约5.08厘米、高约2.11厘米圆角矩形，使用浅蓝色背景，颜色代码为DDE9FF，深蓝色框线，颜色代码为3768B3，字体为Noto Sans CJK SC、约13.2磅、加粗、水平及垂直居中。',
         keys=["2025年", "引入合规与数据", "治理机制"],
         w=5.08, h=2.11, fill="DDE9FF", line="3768B3",
         font="Noto Sans CJK SC", pt=13.2, bold=True, center=True,
         use_strict_rule=True),
    dict(score=3,
         label='"林远航 持股 65% 发起投资"外框：宽约4.45厘米、高约2.08厘米的圆角矩形；使用浅蓝色背景，颜色代码为DDE9FF，深蓝色框线，颜色代码为3768B3，内部字体为Noto Sans CJK SC、约12.3磅、加粗、居中',
         keys=["林远航", "持股 65%", "发起投资"],
         w=4.45, h=2.08, fill="DDE9FF", line="3768B3",
         font="Noto Sans CJK SC", pt=12.3, bold=True, center=True,
         use_strict_inner=True),
    dict(score=3,
         label='"周予安 持股 35% 产业资源"外框：宽约4.45厘米、高约2.08厘米的圆角矩形；使用浅蓝色背景，颜色代码为DDE9FF，深蓝色框线，颜色代码为3768B3，内部字体为Noto Sans CJK SC、约12.3磅、加粗、居中',
         keys=["周予安", "持股 35%", "产业资源"],
         w=4.45, h=2.08, fill="DDE9FF", line="3768B3",
         font="Noto Sans CJK SC", pt=12.3, bold=True, center=True,
         use_strict_inner=True),
    dict(score=3,
         label='"周予安 持股 30% 总经理"外框：宽约4.45厘米、高约2.08厘米的圆角矩形；使用浅蓝色背景，颜色代码为DDE9FF，深蓝色框线，颜色代码为3768B3，内部字体为Noto Sans CJK SC、约12.3磅、加粗、居中',
         keys=["周予安", "持股 30%", "总经理"],
         w=4.45, h=2.08, fill="DDE9FF", line="3768B3",
         font="Noto Sans CJK SC", pt=12.3, bold=True, center=True,
         use_strict_inner=True),
    dict(score=3,
         label='"董事会决策层 战略审批 / 重大合同"外框：宽约4.45厘米、高约2.08厘米的圆角矩形；使用浅蓝色背景，颜色代码为DDE9FF，深蓝色框线，颜色代码为3768B3，内部字体为Noto Sans CJK SC、约12.3磅、加粗、居中',
         keys=["董事会决策层", "战略审批", "重大合同"],
         w=4.45, h=2.08, fill="DDE9FF", line="3768B3",
         font="Noto Sans CJK SC", pt=12.3, bold=True, center=True,
         use_strict_inner=True),
    dict(score=3,
         label='"仓配试点小组 接入订单、库位、车辆数据"外框：宽约4.45厘米、高约2.08厘米的圆角矩形，使用浅粉色背景，颜色代码为F8E3E8，深粉色框线，颜色代码为AD4E66，字体为Noto Sans CJK SC、约12磅、加粗、居中。',
         keys=["仓配试点小组", "接入订单", "库位", "车辆数据"],
         w=4.45, h=2.08, fill="F8E3E8", line="AD4E66",
         font="Noto Sans CJK SC", pt=12.0, bold=True, center=True,
         use_strict_inner=True),
    dict(score=3,
         label='"日常运营链路 订单调度 → 仓储履约 → 异常复盘"外框：宽约4.45厘米、高约2.08厘米的圆角矩形，使用浅粉色背景，颜色代码为F8E3E8，深粉色框线，颜色代码为AD4E66，字体为Noto Sans CJK SC、约12磅、加粗、居中。',
         keys=["日常运营链路", "订单调度", "仓储履约", "异常复盘"],
         w=4.45, h=2.08, fill="F8E3E8", line="AD4E66",
         font="Noto Sans CJK SC", pt=12.0, bold=True, center=True,
         use_strict_inner=True),
    dict(score=3,
         label='"经营执行线 周予安 总经理"外框：宽约4.45厘米、高约2.08厘米的圆角矩形，使用浅粉色背景，颜色代码为F8E3E8，深粉色框线，颜色代码为AD4E66，字体为Noto Sans CJK SC、约12磅、加粗、居中。',
         keys=["经营执行线", "周予安", "总经理"],
         w=4.45, h=2.08, fill="F8E3E8", line="AD4E66",
         font="Noto Sans CJK SC", pt=12.0, bold=True, center=True,
         use_strict_inner=True),
    dict(score=3,
         label='"监督复核线 许知澜 监事"外框：宽约4.45厘米、高约2.08厘米的圆角矩形，使用浅粉色背景，颜色代码为F8E3E8，深粉色框线，颜色代码为AD4E66，字体为Noto Sans CJK SC、约12磅、加粗、居中。',
         keys=["监督复核线", "许知澜", "监事"],
         w=4.45, h=2.08, fill="F8E3E8", line="AD4E66",
         font="Noto Sans CJK SC", pt=12.0, bold=True, center=True,
         use_strict_inner=True),
]


def check_plus_rule_strict_2023(slide, rule):
    """
    严格按细则评估 "2023年 管理负责人调整 职能重新划分" 外框。

    细则原文逐点：
      1. 圆角矩形
      2. 宽约 5.08 厘米
      3. 高约 2.11 厘米
      4. 浅蓝色背景，颜色代码 DDE9FF
      5. 深蓝色框线，颜色代码 3768B3
      6. 字体 Noto Sans CJK SC
      7. 字号约 13.2 磅
      8. 加粗
      9. 水平居中
      10. 垂直居中
    每一条都必须满足，才计 +3。"""
    sh = find_shape_by_text(slide, rule["keys"])
    if sh is None:
        return False, [f'未找到包含 {rule["keys"]} 的形状']

    reasons = []

    # 1) 圆角矩形
    if not is_rounded_rect(sh):
        reasons.append("不是圆角矩形")

    # 3) 宽约 5.08 厘米（"约"取 ±0.3cm）
    w_cm = emu_to_cm(sh.width)
    if not approx(w_cm, rule["w"], tol=0.3):
        reasons.append(f"宽 {w_cm:.2f}cm ≠ 约 {rule['w']}cm")

    # 4) 高约 2.11 厘米
    h_cm = emu_to_cm(sh.height)
    if not approx(h_cm, rule["h"], tol=0.3):
        reasons.append(f"高 {h_cm:.2f}cm ≠ 约 {rule['h']}cm")

    # 5) 浅蓝色背景 DDE9FF（颜色代码明确给出，严格匹配；允许极小渲染误差）
    fill = get_shape_fill_hex(sh)
    if not (fill and str(fill).upper() == rule["fill"].upper()):
        reasons.append(f"背景色 {fill} ≠ {rule['fill']}")

    # 6) 深蓝色框线 3768B3
    line = get_shape_line_hex(sh)
    if not (line and str(line).upper() == rule["line"].upper()):
        reasons.append(f"框线色 {line} ≠ {rule['line']}")

    info = get_font_summary(sh)
    names = info["names"] or set()
    sizes = info["sizes"] or []
    bolds = info["bolds"] or []
    aligns = info["aligns"] or set()

    # 7) 字体 Noto Sans CJK SC（所有 run 都必须是这个字体）
    non_noto = [n for n in names if n != rule["font"]]
    if (rule["font"] not in names) or non_noto:
        reasons.append(f"字体 {sorted(names)} ≠ 仅 {rule['font']}")

    # 8) 字号约 13.2 磅（"约"取 ±0.5pt，所有 run 都要在范围内）
    if not sizes:
        reasons.append("未取到字号")
    else:
        bad_sz = [s for s in sizes if abs(s - rule["pt"]) > 0.5]
        if bad_sz:
            reasons.append(f"字号 {sizes} ≠ 约 {rule['pt']}pt")

    # 9) 加粗（所有 run 都必须 bold=True）
    if not bolds or not all(bolds):
        reasons.append("文字未全部加粗")

    # 10) 水平居中（所有段落对齐都必须是 CENTER）
    if not aligns or aligns != {PP_ALIGN.CENTER}:
        reasons.append(f"未水平居中（aligns={aligns}）")

    # 11) 垂直居中
    if info["vanchor"] != MSO_ANCHOR.MIDDLE:
        reasons.append(f"未垂直居中（vanchor={info['vanchor']}）")

    return (len(reasons) == 0), reasons


def check_plus_rule_strict_inner(slide, rule):
    """
    严格按细则评估"内部框"类（如"林远航 持股 65% 发起投资"）。

    细则原文逐点：
      1. 圆角矩形
      2. 宽约 4.45 厘米
      3. 高约 2.08 厘米
      4. 浅蓝色背景，颜色代码 DDE9FF
      5. 深蓝色框线，颜色代码 3768B3
      6. 内部字体为 Noto Sans CJK SC
      7. 字号约 12.3 磅
      8. 加粗
      9. 居中（水平 + 垂直）
    每一条都必须满足才计 +3。细则未提及的属性（位置、阴影、圆角半径等）一概不约束。"""
    sh = find_shape_by_text(slide, rule["keys"])
    if sh is None:
        return False, [f'未找到包含 {rule["keys"]} 的形状']

    reasons = []

    # 1) 圆角矩形
    if not is_rounded_rect(sh):
        reasons.append("不是圆角矩形")

    # 2) 宽约 4.45 厘米
    w_cm = emu_to_cm(sh.width)
    if not approx(w_cm, rule["w"], tol=0.3):
        reasons.append(f"宽 {w_cm:.2f}cm ≠ 约 {rule['w']}cm")

    # 3) 高约 2.08 厘米
    h_cm = emu_to_cm(sh.height)
    if not approx(h_cm, rule["h"], tol=0.3):
        reasons.append(f"高 {h_cm:.2f}cm ≠ 约 {rule['h']}cm")

    # 4) 浅蓝色背景 DDE9FF（颜色代码明确给出，严格匹配）
    fill = get_shape_fill_hex(sh)
    if not (fill and str(fill).upper() == rule["fill"].upper()):
        reasons.append(f"背景色 {fill} ≠ {rule['fill']}")

    # 5) 深蓝色框线 3768B3
    line = get_shape_line_hex(sh)
    if not (line and str(line).upper() == rule["line"].upper()):
        reasons.append(f"框线色 {line} ≠ {rule['line']}")

    info = get_font_summary(sh)
    names = info["names"] or set()
    sizes = info["sizes"] or []
    bolds = info["bolds"] or []
    aligns = info["aligns"] or set()

    # 6) 内部字体 Noto Sans CJK SC（所有 run 都必须是该字体）
    non_noto = [n for n in names if n != rule["font"]]
    if (rule["font"] not in names) or non_noto:
        reasons.append(f"字体 {sorted(names)} ≠ 仅 {rule['font']}")

    # 7) 字号约 12.3 磅（每个 run 都在 ±0.5pt 内）
    if not sizes:
        reasons.append("未取到字号")
    else:
        bad_sz = [s for s in sizes if abs(s - rule["pt"]) > 0.5]
        if bad_sz:
            reasons.append(f"字号 {sizes} ≠ 约 {rule['pt']}pt")

    # 8) 加粗（所有 run）
    if not bolds or not all(bolds):
        reasons.append("文字未全部加粗")

    # 9) 居中（水平 + 垂直；细则只说"居中"，按业内常规两者都满足才算"居中"）
    if not aligns or aligns != {PP_ALIGN.CENTER}:
        reasons.append(f"未水平居中（aligns={aligns}）")
    if info["vanchor"] != MSO_ANCHOR.MIDDLE:
        reasons.append(f"未垂直居中（vanchor={info['vanchor']}）")

    return (len(reasons) == 0), reasons


def check_plus_rule(slide, rule):
    """返回 (matched: bool, reasons: list[str])"""
    # 严格按细则评估的项走专用函数
    if rule.get("use_strict_rule"):
        return check_plus_rule_strict_2023(slide, rule)
    if rule.get("use_strict_inner"):
        return check_plus_rule_strict_inner(slide, rule)

    sh = find_shape_by_text(slide, rule["keys"])
    if sh is None:
        return False, [f'未找到包含 {rule["keys"]} 的形状']
    reasons = []
    # 形状必须是圆角矩形
    if not is_rounded_rect(sh):
        reasons.append("不是圆角矩形")
    # 尺寸
    w_cm = emu_to_cm(sh.width); h_cm = emu_to_cm(sh.height)
    if not approx(w_cm, rule["w"], tol=0.3):
        reasons.append(f"宽 {w_cm:.2f}cm ≠ {rule['w']}cm")
    if not approx(h_cm, rule["h"], tol=0.3):
        reasons.append(f"高 {h_cm:.2f}cm ≠ {rule['h']}cm")
    # 填充
    fill = get_shape_fill_hex(sh)
    if not color_close(fill, rule["fill"], tol=20):
        reasons.append(f"填充 {fill} ≠ {rule['fill']}")
    # 线
    line = get_shape_line_hex(sh)
    if not color_close(line, rule["line"], tol=24):
        reasons.append(f"线色 {line} ≠ {rule['line']}")
    # 字体
    info = get_font_summary(sh)
    if rule["font"] not in info["names"]:
        reasons.append(f"字体 {info['names']} 不含 {rule['font']}")
    if info["sizes"]:
        avg_pt = sum(info["sizes"]) / len(info["sizes"])
        if abs(avg_pt - rule["pt"]) > 1.5:
            reasons.append(f"字号均值 {avg_pt:.1f}pt ≠ {rule['pt']}pt")
    else:
        reasons.append("未取到字号")
    if rule["bold"] and info["bolds"] and not all(info["bolds"]):
        reasons.append("文字未全部加粗")
    # 居中
    if rule["center"]:
        if not alignment_ok(info):
            reasons.append(f"未水平/垂直居中 aligns={info['aligns']} vanchor={info['vanchor']}")
    return (len(reasons) == 0), reasons


# ---------- 维度2 扣分项 ----------

# （原三项扣分项已按需求删除）


# ---------- 统一入口 ----------

SCRIPT_ID = "069"


def _locate_pptx(dir_path: str):
    """在给定目录下定位待评估的 .pptx 文件。返回绝对路径或 None。"""
    if not os.path.isdir(dir_path):
        return None
    candidates = [
        f for f in os.listdir(dir_path)
        if f.lower().endswith(".pptx") and not f.startswith("~$")
    ]
    if not candidates:
        return None
    # 若存在多个，按文件名排序取第一个（结果稳定）
    candidates.sort()
    return os.path.join(dir_path, candidates[0])


def evaluate(dir_path: str) -> dict:
    """统一入口：接收脚本所在目录的路径，脚本自行在该目录中定位并打开被评估文档。

    返回结构见 §2.2：
      {
        "id", "file_name", "status", "error",
        "dim1_pass", "dim1_reason",
        "dim2_items": [{"rule","max_delta","delta","hit","detail"}, ...],
        "total_score", "max_score"
      }
    """
    max_score = sum(r["score"] for r in PLUS_RULES)
    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": max_score,
    }

    try:
        pptx_path = _locate_pptx(dir_path)
        if pptx_path is None:
            result["status"] = "error"
            result["error"] = f"目录中未找到 .pptx 文件：{dir_path}"
            return result
        result["file_name"] = os.path.basename(pptx_path)

        # 维度 1
        d1_pass, d1_details = check_dimension_1(pptx_path)
        result["dim1_pass"] = bool(d1_pass)
        if not d1_pass:
            result["dim1_reason"] = "；".join(d1_details) if d1_details else ""
            result["total_score"] = 0
            return result

        # 维度 2
        prs = Presentation(pptx_path)
        slide = prs.slides[0]

        total = 0
        items = []

        # 加分项（命中和未命中都列出）
        for rule in PLUS_RULES:
            ok, _ = check_plus_rule(slide, rule)
            delta = rule["score"] if ok else 0
            total += delta
            items.append({
                "rule": rule["label"],
                "max_delta": rule["score"],
                "delta": delta,
                "hit": bool(ok),
                "detail": "",
            })

        result["dim2_items"] = items
        result["total_score"] = total
        return result

    except Exception as e:
        result["status"] = "error"
        result["error"] = f"{type(e).__name__}: {e}"
        return result


if __name__ == "__main__":
    # 本地调试入口：默认使用当前脚本所在目录；也可通过命令行覆盖
    import json
    target_dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(target_dir), ensure_ascii=False, indent=2, default=str))
