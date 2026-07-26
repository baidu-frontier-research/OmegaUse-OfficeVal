# -*- coding: utf-8 -*-
"""
自动评估 PPT 文件：农业县低空经济发展方案_English.pptx

评估流程：
  维度1（可用与可修改性）—— 任一不满足直接判为 0 分，不再检查维度2。
  维度2（完成度评分细则）—— 累计每条命中的细则分值（加 / 扣）。

加分细则：必须 *全部* 满足该条所列的所有点，才计入加分。
扣分细则：只要满足该条所列的 *任意一点*，即计入扣分。

最终输出：命中的细则明细 + 最终得分。
"""

import json
import os
import re
import sys

from pptx import Presentation

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ----------------------------- 配置 -----------------------------

SCRIPT_ID = "074"
# 期望的目标文档名（若在传入目录里找不到该名，会退化为该目录下任意 .pptx 文件）
EXPECTED_DOC_NAME = "农业县低空经济发展方案_English.pptx"

EXPECTED_PAGES = 46
TITLE_SLIDES = {3, 6, 14, 17, 19, 23, 42}     # 章节标题页
COVER_SLIDE = 1
TOC_SLIDE = 2
THANK_SLIDE = 46
NON_CONTENT_SLIDES = {1, 2, 3, 6, 14, 17, 19, 23, 42, 46}

# 字号容差（pt）。PPT 中字号常存在 ±0.5pt 的舍入。
SZ_TOL = 0.5

# ----------------------------- 工具函数 -----------------------------


def iter_runs(slide):
    """yield (shape, paragraph, run, text, raw_size_pt_or_None, bold_or_None, font_name)"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                txt = run.text or ""
                if not txt.strip():
                    continue
                sz = run.font.size.pt if run.font.size else None
                bold = run.font.bold
                fname = run.font.name
                yield shape, para, run, txt, sz, bold, fname


def slide_full_text(slide):
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            parts.append(shape.text_frame.text)
    return "\n".join(parts)


def about(a, b, tol=SZ_TOL):
    return a is not None and abs(a - b) <= tol


def find_page_footer_run(slide, idx):
    """页码行：例如 '01  Plan Cover  ·  Start  →  Next: Contents'。
    通过以两位/三位页码数字开头并包含 '·' 或 '→' 来识别。"""
    prefix = f"{idx:02d}"
    for shape, para, run, txt, sz, bold, fname in iter_runs(slide):
        t = txt.strip()
        if t.startswith(prefix) and ("·" in t or "→" in t or "Previous" in t or "Next" in t or "Wrap" in t):
            return run, sz, t
    return None, None, None


# ----------------------------- 维度1 -----------------------------


def check_dim1(path: str):
    """维度 1 硬体检，返回 (pass, reasons)。

    规则简化（原细则 "交付文件为.pptx格式，文件可正常打开、编辑和播放" 已简化为
    "交付文件为.pptx格式，文件可正常打开"）：
      • 后缀为 .pptx；
      • 文件存在，且可被 python-pptx 正常打开。
    之前的 "文档保护 / 形状可编辑锁 / PowerPoint COM 播放性" 三项校验已弃用。
    页数 = EXPECTED_PAGES 这一独立于本条的硬性要求保留。
    """
    reasons = []
    ok = True

    # 1) 格式 .pptx，能正常打开
    ext = os.path.splitext(path)[1].lower()
    if ext != ".pptx":
        reasons.append(f"文件扩展名不是 .pptx：{ext}")
        return False, reasons

    if not os.path.isfile(path):
        reasons.append("文件不存在")
        return False, reasons

    # 2) 能否被 python-pptx 正常打开（作为 pptx 包结构合法的最小充分条件）
    try:
        prs = Presentation(path)
    except Exception as e:
        reasons.append(f"文件无法正常打开：{e}")
        return False, reasons

    # 3) 页数 = EXPECTED_PAGES
    n = len(prs.slides)
    if n != EXPECTED_PAGES:
        reasons.append(f"页数为 {n}，不等于 {EXPECTED_PAGES}")
        ok = False

    if ok:
        reasons.append(f"✓ 文件为 {ext}，可正常打开；页数 {n}")

    return ok, reasons


# ----------------------------- 维度2：判定函数 -----------------------------


# ===== 加分项 =====


def rule_p1_cover_title_translated(prs):
    """+1 第1页主标题文本框：中文"青岚县农业低空应用实施方案"
        翻译为"Qinglan County Agricultural Low-Altitude Application Implementation Plan"
        或语义一致英文。

    严格对齐细则原文，只判定细则要求的事——
    第 1 页的主标题文本框，文本是"指定英文"或"语义一致英文"。
    其它（中文原文是否还在、是否加粗、字号几何…）不在本条细则内，
    不在此处加以约束。
    """
    slide = prs.slides[COVER_SLIDE - 1]                                # 第 1 页
    expected_en = "Qinglan County Agricultural Low-Altitude Application Implementation Plan"

    # 主标题文本框：取该页字号最大的文本框（PPT 中"主标题"的通常定位方式）
    main_title_text = None
    main_title_size = -1
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        shape_text = shape.text_frame.text.strip()
        if not shape_text:
            continue
        max_sz_in_shape = -1
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not (run.text or "").strip():
                    continue
                sz = run.font.size.pt if run.font.size else 0
                if sz > max_sz_in_shape:
                    max_sz_in_shape = sz
        if max_sz_in_shape > main_title_size:
            main_title_size = max_sz_in_shape
            main_title_text = shape_text

    if not main_title_text:
        return False, "第1页未找到主标题文本框"

    # 指定英文（不区分大小写包含匹配）
    is_exact = expected_en.lower() in main_title_text.lower()
    # 语义一致英文：用与中文原文逐词对应的英文语义判断
    # "青岚县" → Qinglan(County) / "农业" → agricultural / "低空" → low altitude /
    # "应用" → application(s) / "实施方案" → implementation plan
    low = main_title_text.lower()
    has_place    = "qinglan" in low
    has_agri     = "agricultur" in low                                  # agricultural / agriculture
    has_lowalt   = ("low-altitude" in low) or ("low altitude" in low)
    has_app      = "application" in low                                 # application / applications
    has_implplan = ("implementation plan" in low) or ("implementation" in low and "plan" in low)
    is_semantic  = has_place and has_agri and has_lowalt and has_app and has_implplan

    hit = is_exact or is_semantic
    detail = (
        f"主标题文本框='{main_title_text[:80]}' | "
        f"指定英文={is_exact}, 语义一致英文={is_semantic}"
    )
    return hit, detail


def rule_p2_toc_all_english(prs):
    """+1 第2页目录页文本：目录标题、PART 01至PART 07、各部分名称和说明全部为英文。

    严格对齐细则的每一个点：
      (a) 第 2 页；
      (b) 目录标题——存在；且为英文（含合理英文词汇 + 不含非英文主要文字）；
          目录标题不硬编码为"Contents"——先按 placeholder 定位目录标题对象，
          失败再匹配已知英文目录标题词汇集合。
      (c) PART 01 至 PART 07——七个标记齐全；且各自所在文本为英文；
      (d) 各部分名称（PART 标题的后半段，如 "Project Background"）为英文；
      (e) 各部分说明（紧跟在 PART 标题之后的描述文字）为英文。
    """
    slide = prs.slides[TOC_SLIDE - 1]                                  # (a) 第 2 页

    # ---- 英文判定 ----
    # 非英文主要文字区间：CJK / 假名 / 韩文 / 西里尔 / 阿拉伯 / 希伯来 / 天城文 / 泰文 / 希腊
    NON_EN_RANGES = (
        ("一", "鿿"), ("㐀", "䶿"),
        ("぀", "ゟ"), ("゠", "ヿ"),
        ("가", "힯"),
        ("Ѐ", "ӿ"),
        ("؀", "ۿ"), ("֐", "׿"),
        ("ऀ", "ॿ"), ("฀", "๿"),
        ("Ͱ", "Ͽ"),
    )
    # 强指示非英文的拉丁扩展字符（法/德/西/葡等重音字母），若出现直接判非英文
    NON_EN_LATIN = set(
        "àâäæçèéêëîïôœùûüÿñÀÂÄÆÇÈÉÊËÎÏÔŒÙÛÜŸÑßẞ" +
        "áíóúÁÍÓÚãõÃÕ"
    )
    EN_WORD_RE = re.compile(r"[A-Za-z][A-Za-z'\-]{1,}")

    def has_non_english_script(s: str) -> bool:
        for c in s:
            if c in NON_EN_LATIN:
                return True
            for lo, hi in NON_EN_RANGES:
                if lo <= c <= hi:
                    return True
        return False

    def is_english_text(s: str) -> bool:
        """英文判定：不含非英文主要文字；至少一个长度≥2的英文单词；
        ASCII 字母 / 所有字母 ≥ 0.9（排除拼音/乱码/纯数字）。"""
        if not s or not s.strip():
            return False
        if has_non_english_script(s):
            return False
        if not EN_WORD_RE.search(s):
            return False
        ascii_letters = sum(1 for c in s if ("A" <= c <= "Z") or ("a" <= c <= "z"))
        all_letters = sum(1 for c in s if c.isalpha())
        if all_letters == 0 or ascii_letters / all_letters < 0.9:
            return False
        return True

    # 已知英文目录标题词汇集合（用于 placeholder 定位失败时的兜底匹配）
    KNOWN_EN_TOC_TITLES = {
        "contents", "table of contents", "toc", "index", "outline",
        "agenda", "directory", "chapters", "sections",
    }

    # 把该页所有"非空 run"按出现顺序收集起来，便于做"PART 标题 → 紧随说明"的配对
    runs_seq = []  # [(text, shape_name)]
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                t = (run.text or "").strip()
                if t:
                    runs_seq.append((t, shape.name))

    missing = []   # 未找到的点
    bad_lang = []  # 找到了但不是英文

    # (b) 目录标题：先按目录标题对象（TITLE / CENTER_TITLE placeholder）定位；
    #     失败再退化为"页面上半部字号最大的非 PART 文本框"或"匹配已知英文目录标题词"。
    title_shape = None
    for sh in slide.shapes:
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        try:
            if sh.is_placeholder and sh.placeholder_format is not None:
                pht = sh.placeholder_format.type
                # 13 = TITLE, 15 = CENTER_TITLE（python-pptx 的枚举整数值）
                if pht is not None and int(pht) in (13, 15):
                    title_shape = sh
                    break
        except Exception:
            pass

    title_text = None
    if title_shape is not None:
        title_text = title_shape.text_frame.text.strip()
    else:
        # 退化 1：页面上半部字号最大的非 "PART 0X" 且非页脚的文本框
        best_sh, best_sz = None, -1
        page_top_half = (prs.slide_height or 0) / 2
        for sh in slide.shapes:
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            t = sh.text_frame.text.strip()
            if t.startswith("PART 0"):
                continue
            if "·" in t and ("→" in t or "Previous" in t or "Next" in t):
                continue
            if (sh.top or 0) > page_top_half:
                continue
            mx = -1
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        sz = run.font.size.pt if run.font.size else 0
                        if sz > mx:
                            mx = sz
            if mx > best_sz:
                best_sz, best_sh = mx, sh
        if best_sh is not None:
            title_text = best_sh.text_frame.text.strip()
        else:
            # 退化 2：匹配已知英文目录标题词
            for t, _ in runs_seq:
                if t.strip().lower() in KNOWN_EN_TOC_TITLES:
                    title_text = t.strip()
                    break

    if not title_text:
        missing.append("目录标题")
    else:
        # 目录标题必须为英文；额外要求它"看起来像目录标题"——是已知英文目录标题词
        # 或包含合理英文词汇（不再硬编码为 Contents）
        if not is_english_text(title_text):
            bad_lang.append(f"目录标题非英文：'{title_text[:40]}'")

    # (c)(d)(e) PART 01 至 PART 07：标题 + 紧随其后的"说明"配对
    for n in range(1, 8):
        tag = f"PART 0{n}"
        idx_part = None
        for i, (t, _) in enumerate(runs_seq):
            if t.startswith(tag):
                idx_part = i
                break
        if idx_part is None:
            missing.append(f"{tag} 标题")
            continue

        # PART 标题整段（含 "PART 0X" 与部分名称）必须为英文
        part_title_text = runs_seq[idx_part][0]
        # "PART 0X" 前缀本身是英文，去掉后剩余部分才是"部分名称"，用它做主要英文判定
        part_name = part_title_text[len(tag):].strip()
        if has_non_english_script(part_title_text):
            bad_lang.append(f"{tag} 标题含非英文文字：'{part_title_text}'")
        if not part_name:
            missing.append(f"{tag} 各部分名称")
        elif not is_english_text(part_name):
            bad_lang.append(f"{tag} 各部分名称非英文：'{part_name}'")

        # 各部分说明：紧跟在该 PART 标题之后、且不是下一个 PART 标题 / 页脚的 run
        desc = None
        for j in range(idx_part + 1, len(runs_seq)):
            nxt = runs_seq[j][0]
            if nxt.startswith("PART 0"):
                break
            # 跳过页脚行
            if "·" in nxt and ("→" in nxt or "Previous" in nxt or "Next" in nxt):
                continue
            desc = nxt
            break
        if desc is None:
            missing.append(f"{tag} 各部分说明")
        elif not is_english_text(desc):
            bad_lang.append(f"{tag} 各部分说明非英文：'{desc}'")

    hit = (not missing) and (not bad_lang)
    parts = []
    parts.append("目录标题=" + (f"'{title_text[:30]}'" if title_text else "缺失"))
    parts.append(f"PART01-07 标题齐全={all(any(t.startswith(f'PART 0{n}') for t,_ in runs_seq) for n in range(1,8))}")
    if missing:
        parts.append("缺失:" + "/".join(missing))
    if bad_lang:
        parts.append("非英文:" + "/".join(bad_lang))
    detail = "目录页英文化检查：" + "； ".join(parts)
    return hit, detail


def rule_p3_to_p46_all_english(prs):
    """+3 第3页至第46页内容文本全部为英文。

    严格对齐细则的每一个点：
      (a) 范围：第 3 页 至 第 46 页（含两端）；
      (b) 内容文本：这些页面上所有文本框中的文本；
      (c) 全部为英文：非英文脚本(CJK/假名/韩文/西里尔/阿拉伯/希伯来/天城文/泰文/希腊)
          任一命中即非英文；同时排除"不含中文但也不是英文"的情况——
          要求文本中至少含合理英文单词、且英文字符/字母比例达到阈值；
          若安装了 langdetect 且文本较长(>=30字母)，还会做语言识别兜底。
    """
    # ---- 非英文脚本 / 强指示非英文的拉丁扩展字符 ----
    NON_EN_RANGES = (
        ("一", "鿿"), ("㐀", "䶿"),      # CJK 中日韩统一表意
        ("぀", "ゟ"), ("゠", "ヿ"),      # 平假 / 片假
        ("가", "힯"),                    # 韩文
        ("Ѐ", "ӿ"),                     # 西里尔
        ("؀", "ۿ"), ("֐", "׿"),        # 阿拉伯 / 希伯来
        ("ऀ", "ॿ"), ("฀", "๿"),        # 天城文 / 泰文
        ("Ͱ", "Ͽ"),                     # 希腊
    )
    NON_EN_LATIN = set(
        "àâäæçèéêëîïôœùûüÿñÀÂÄÆÇÈÉÊËÎÏÔŒÙÛÜŸÑßẞ" +
        "áíóúÁÍÓÚãõÃÕ"
    )
    EN_WORD_RE = re.compile(r"[A-Za-z][A-Za-z'\-]{1,}")

    # 可选：语言识别（有则更稳，无则跳过）
    # 用 importlib 动态加载，避免静态分析器在未安装 langdetect 时报"无法解析导入"
    import importlib
    _detect_fn = None
    try:
        _ld = importlib.import_module("langdetect")
        _detect_fn = getattr(_ld, "detect", None)
        _factory = getattr(_ld, "DetectorFactory", None)
        if _factory is not None:
            try:
                _factory.seed = 0
            except Exception:
                pass
    except Exception:
        _detect_fn = None
    _has_langdetect = _detect_fn is not None

    def find_non_en_script(s: str):
        for c in s:
            if c in NON_EN_LATIN:
                return c
            for lo, hi in NON_EN_RANGES:
                if lo <= c <= hi:
                    return c
        return None

    def classify(text: str):
        """返回 (is_english, reason)。reason 为非英文时的简短原因。"""
        s = text or ""
        stripped = s.strip()
        if not stripped:
            return True, ""   # 空文本不构成"非英文"

        # 1) 非英文脚本 / 强非英拉丁扩展直接判非英文
        bad_ch = find_non_en_script(s)
        if bad_ch is not None:
            return False, f"含非英文字符'{bad_ch}'"

        # 2) 英文单词存在性（长度≥2 的字母序列）
        if not EN_WORD_RE.search(s):
            # 无字母（如纯数字/符号）不判为"非英文"内容，跳过
            if not any(c.isalpha() for c in s):
                return True, ""
            return False, "无可识别英文单词"

        # 3) 英文字母比例
        ascii_letters = sum(1 for c in s if ("A" <= c <= "Z") or ("a" <= c <= "z"))
        all_letters = sum(1 for c in s if c.isalpha())
        if all_letters and ascii_letters / all_letters < 0.9:
            return False, f"英文字母比例过低({ascii_letters}/{all_letters})"

        # 4) 单词占比：英文单词字符 / 全部字母 ≥ 0.6，避免"一串重音+个别英文词"
        word_chars = sum(len(m.group()) for m in EN_WORD_RE.finditer(s))
        if all_letters and word_chars / all_letters < 0.6:
            return False, f"英文单词字符占比过低({word_chars}/{all_letters})"

        # 5) 语言识别兜底：对足够长的文本调用 langdetect
        if _has_langdetect and _detect_fn is not None and all_letters >= 30:
            try:
                lang = _detect_fn(s)
                if lang != "en":
                    return False, f"语言识别为 '{lang}'"
            except Exception:
                pass

        return True, ""

    bad = []  # [(page, sample_text, reason)]
    for page_no in range(3, EXPECTED_PAGES + 1):                       # (a) 第3-46页
        slide = prs.slides[page_no - 1]
        for shape in slide.shapes:                                     # (b) 内容文本
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if not text or not text.strip():
                continue
            is_en, reason = classify(text)
            if not is_en:
                sample = "".join(c for c in text if not c.isspace())[:40]
                bad.append((page_no, sample, reason))
                break  # 该页已记录一次即可

    hit = len(bad) == 0
    if hit:
        tail = "（含 langdetect 兜底）" if _has_langdetect else "（未安装 langdetect，未做语言识别兜底）"
        detail = f"第3-46页所有内容文本均判定为英文{tail}"
    else:
        show = bad[:8]
        more = f"…等共 {len(bad)} 页" if len(bad) > 8 else ""
        detail = "存在非英文内容的页：" + ", ".join(
            f"第{p}页[{r}]({s})" for p, s, r in show
        ) + (" " + more if more else "")
    return hit, detail


# ===== 扣分项 =====


def rule_m1_cover_title_font(prs):
    """-1 首页"Qinglan County Agricultural Low-Altitude Application Implementation Plan"
        字体不是 32 磅加粗。

    严格对齐细则的每一个点：
      (a) 首页 —— 第 1 页；
      (b) 目标文本 —— "Qinglan County Agricultural Low-Altitude Application Implementation Plan"；
      (c) 字号是 32 磅；
      (d) 加粗。
      其中 (c)(d) 任一不满足即扣分（"不是 32 磅加粗"是对"32 磅 且 加粗"的否定）。
    """
    slide = prs.slides[COVER_SLIDE - 1]                                # (a) 首页
    expected_text = "Qinglan County Agricultural Low-Altitude Application Implementation Plan"

    # (b) 在首页查找包含该指定英文文本的 run
    target_size = None
    target_bold = None
    found = False
    for _, _, _, txt, sz, bold, _ in iter_runs(slide):
        if expected_text.lower() in txt.lower():
            target_size = sz
            target_bold = bold
            found = True
            break

    if not found:
        # 细则只在"首页存在该指定文本"的前提下检查字号/加粗。
        # 文本不存在时，本条不命中扣分。
        return False, f"首页未找到文本'{expected_text}'，本条扣分不适用"

    # (c)(d) 32 磅 且 加粗
    size_ok = about(target_size, 32.0)
    bold_ok = bool(target_bold)
    is_32pt_bold = size_ok and bold_ok
    hit = not is_32pt_bold                                              # 不是 32 磅加粗 → 命中扣分
    detail = (
        f"首页'{expected_text[:50]}...'实测：字号={target_size}磅(需32磅), "
        f"加粗={target_bold}(需True)"
    )
    return hit, detail


def rule_m1_cover_subtitle_font(prs):
    """-1 首页"Qinglan County Agricultural Low-Altitude Application Implementation Plan"
        下方文本字体不是 13 磅。

    严格对齐细则的每一个点：
      (a) 首页 —— 第 1 页；
      (b) 锚点 —— 文本"Qinglan County Agricultural Low-Altitude Application Implementation Plan"
          所在的文本框（即首页主标题文本框）；
      (c) "下方文本" —— 在锚点文本框正下方（top 大于锚点的 top + height/2，且水平方向有重叠）
          的最近一个文本框中的文本；
      (d) 字体是 13 磅。
      不是 13 磅则命中扣分。
    """
    slide = prs.slides[COVER_SLIDE - 1]                                # (a) 首页
    expected_text = "Qinglan County Agricultural Low-Altitude Application Implementation Plan"

    # (b) 找锚点文本框
    anchor = None
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if expected_text.lower() in shape.text_frame.text.lower():
            anchor = shape
            break
    if anchor is None:
        return False, f"首页未找到锚点文本'{expected_text}'，本条扣分不适用"

    a_top    = anchor.top or 0
    a_left   = anchor.left or 0
    a_right  = a_left + (anchor.width or 0)
    a_bottom = a_top  + (anchor.height or 0)

    # (c) 在锚点下方，水平方向与锚点有重叠的最近一个文本框
    below = None  # (top, shape)
    for shape in slide.shapes:
        if shape is anchor or not shape.has_text_frame:
            continue
        if not shape.text_frame.text.strip():
            continue
        s_top   = shape.top or 0
        s_left  = shape.left or 0
        s_right = s_left + (shape.width or 0)
        # "下方"：顶部在锚点底部之下（允许小幅交叠）
        if s_top < a_bottom - (anchor.height or 0) * 0.3:
            continue
        # 水平重叠
        if s_right <= a_left or s_left >= a_right:
            continue
        if below is None or s_top < below[0]:
            below = (s_top, shape)

    if below is None:
        return False, "首页未找到主标题下方文本，本条扣分不适用"

    sub_shape = below[1]
    # (d) 取下方文本框内所有非空 run 的字号集合
    sizes = []
    for para in sub_shape.text_frame.paragraphs:
        for run in para.runs:
            if (run.text or "").strip():
                sz = run.font.size.pt if run.font.size else None
                sizes.append(sz)

    is_13pt = (len(sizes) > 0) and all(about(sz, 13.0) for sz in sizes)
    hit = not is_13pt                                                  # 不是 13 磅 → 命中扣分
    sub_text_preview = sub_shape.text_frame.text.strip().replace("\n", " ")[:60]
    detail = f"主标题下方文本='{sub_text_preview}...' 字号集合={sizes}（需 13 磅）"
    return hit, detail


def rule_m5_title_slides_font(prs):
    """-5 PPT标题页（PPT第3、6、14、17、19、23、42页）标题文本字号不是34磅，
        标题下方文本字号不是14磅。

    严格对齐细则的每一个点：
      (a) 范围 —— 第 3、6、14、17、19、23、42 页（仅这 7 页）；
      (b) 标题文本 —— 这些页的"标题"文本框；按以下顺序定位：
          1) 占位符类型为 TITLE(13) / CENTER_TITLE(15)；
          2) 与该章节已知英文标题（PART 01~07 的部分名称）匹配；
          3) 排除装饰性大数字（纯数字文本框，如 "01"）与页脚后，
             再退化为字号最大、位置靠上的文本框。
      (c) 标题字号是 34 磅；
      (d) 标题下方文本 —— 副标题；按以下顺序定位：
          1) 占位符类型为 SUBTITLE(4) / BODY(2)；
          2) 与标题文本框同一形状内、位于标题段落之后的其它段落；
          3) 位于标题文本框下方、水平方向重叠的最近文本框。
      (e) 该下方文本字号是 14 磅。
      "标题非 34 磅" 或 "下方文本非 14 磅"，任一不满足 → 命中扣分。
    """
    # 各标题页对应的已知英文章节名（与 PART 01~07 顺序一一对应）
    KNOWN_TITLES_BY_PAGE = {
        3:  ["project background"],
        6:  ["current status", "current status and challenges", "status", "challenges"],
        14: ["overall design", "overall plan", "design", "plan"],
        17: ["implementation", "implementation plan", "implementation path"],
        19: ["safety", "safety and compliance", "compliance"],
        23: ["benefits", "benefit analysis", "expected benefits"],
        42: ["roadmap", "future outlook", "outlook", "next steps"],
    }
    # python-pptx PP_PLACEHOLDER 枚举整数值
    PH_TITLE, PH_CENTER_TITLE = 13, 15
    PH_SUBTITLE, PH_BODY = 4, 2

    def is_footer_text(t: str) -> bool:
        return "·" in t and ("→" in t or "Previous" in t or "Next" in t or "Wrap" in t)

    def para_text(para) -> str:
        return "".join((r.text or "") for r in para.runs).strip()

    def para_sizes(para):
        return [r.font.size.pt if r.font.size else None
                for r in para.runs if (r.text or "").strip()]

    def shape_sizes(sh):
        out = []
        for para in sh.text_frame.paragraphs:
            out.extend(para_sizes(para))
        return out

    def max_run_size(sh):
        sizes = [s for s in shape_sizes(sh) if s is not None]
        return max(sizes) if sizes else -1

    def ph_type(sh):
        try:
            if sh.is_placeholder and sh.placeholder_format is not None:
                pht = sh.placeholder_format.type
                return int(pht) if pht is not None else None
        except Exception:
            return None
        return None

    bad = []
    for idx in sorted(TITLE_SLIDES):                                   # (a) 7 个标题页
        slide = prs.slides[idx - 1]

        # 收集该页非空、非页脚文本框
        text_shapes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            if not txt or is_footer_text(txt):
                continue
            text_shapes.append(shape)

        # ---- (b) 标题定位 ----
        title_shape = None

        # 1) 占位符类型
        for sh in text_shapes:
            if ph_type(sh) in (PH_TITLE, PH_CENTER_TITLE):
                title_shape = sh
                break

        # 2) 已知英文章节名匹配（大小写不敏感、去空白/标点）
        if title_shape is None:
            def norm(s: str) -> str:
                return re.sub(r"[^a-z0-9]+", " ", s.lower()).strip()
            candidates = [norm(t) for t in KNOWN_TITLES_BY_PAGE.get(idx, [])]
            for sh in text_shapes:
                t = norm(sh.text_frame.text)
                if any(c and c in t for c in candidates):
                    title_shape = sh
                    break

        # 3) 排除装饰性纯数字后，取字号最大且靠上
        if title_shape is None:
            for sh in text_shapes:
                t = sh.text_frame.text.strip()
                if t.isdigit():
                    continue
                if title_shape is None:
                    title_shape = sh
                    continue
                if max_run_size(sh) > max_run_size(title_shape):
                    title_shape = sh
                elif max_run_size(sh) == max_run_size(title_shape) \
                        and (sh.top or 0) < (title_shape.top or 0):
                    title_shape = sh

        if title_shape is None:
            bad.append(f"第{idx}页：未找到标题文本框")
            continue

        # ---- 找到标题段落索引（用于"同一文本框内后续段落即副标题"的场景）----
        title_paras = list(title_shape.text_frame.paragraphs)
        # 标题所在段落：取"最大字号所在段落"（若整框统一字号则为第一非空段落）
        title_para_idx = None
        best_max = -1
        for pi, p in enumerate(title_paras):
            sizes = [s for s in para_sizes(p) if s is not None]
            if not sizes:
                continue
            mx = max(sizes)
            if mx > best_max:
                best_max, title_para_idx = mx, pi
        if title_para_idx is None:
            # 全部无字号信息时退化为第一非空段落
            for pi, p in enumerate(title_paras):
                if para_text(p):
                    title_para_idx = pi
                    break

        # ---- (c) 标题字号 = 34 磅（取标题段落所在 run 字号集合）----
        if title_para_idx is not None:
            title_sizes = para_sizes(title_paras[title_para_idx])
        else:
            title_sizes = shape_sizes(title_shape)
        title_ok = len(title_sizes) > 0 and all(about(sz, 34.0) for sz in title_sizes)

        # ---- (d) 标题下方文本定位 ----
        sub_sizes = []
        sub_source = None  # "placeholder" / "same-shape" / "below-shape"

        # 1) SUBTITLE / BODY 占位符
        for sh in text_shapes:
            if sh is title_shape:
                continue
            if ph_type(sh) in (PH_SUBTITLE, PH_BODY):
                sub_sizes = shape_sizes(sh)
                sub_source = "placeholder"
                below_ref = sh
                break
        else:
            below_ref = None

        # 2) 同一文本框内、标题段落之后的其它非空段落
        if not sub_sizes and title_para_idx is not None:
            same_shape_sizes = []
            for p in title_paras[title_para_idx + 1:]:
                if para_text(p):
                    same_shape_sizes.extend(para_sizes(p))
            if same_shape_sizes:
                sub_sizes = same_shape_sizes
                sub_source = "same-shape"
                below_ref = title_shape

        # 3) 位于标题文本框下方、水平方向重叠的最近文本框
        if not sub_sizes:
            t_top    = title_shape.top or 0
            t_left   = title_shape.left or 0
            t_right  = t_left + (title_shape.width or 0)
            t_bottom = t_top + (title_shape.height or 0)
            below = None
            for sh in text_shapes:
                if sh is title_shape:
                    continue
                # 跳过装饰性纯数字
                if sh.text_frame.text.strip().isdigit():
                    continue
                s_top   = sh.top or 0
                s_left  = sh.left or 0
                s_right = s_left + (sh.width or 0)
                if s_top < t_bottom - (title_shape.height or 0) * 0.3:
                    continue
                if s_right <= t_left or s_left >= t_right:
                    continue
                if below is None or s_top < (below.top or 0):
                    below = sh
            if below is not None:
                sub_sizes = shape_sizes(below)
                sub_source = "below-shape"
                below_ref = below

        # (e) 下方文本字号 = 14 磅
        sub_ok = len(sub_sizes) > 0 and all(about(sz, 14.0) for sz in sub_sizes)

        if not (title_ok and sub_ok):
            src = sub_source or "未找到"
            title_preview = title_shape.text_frame.text.strip().replace("\n", " ")[:30]
            bad.append(
                f"第{idx}页：标题文本框='{title_preview}' "
                f"字号={title_sizes}(需34磅); "
                f"下方文本[来源={src}] 字号={sub_sizes}(需14磅)"
            )

    hit = len(bad) > 0
    detail = "标题页字号检查：\n    " + "\n    ".join(bad) if bad else "7 个标题页字号均合规"
    return hit, detail


def rule_m5_page_footer_font(prs):
    """-5 页码处"01  Plan Cover  ·  Start  →  Next: Contents"等文本字号不是7.8磅。

    严格对齐细则的每一个点：
      (a) 范围 —— 整本 PPT 的"页码处"文本；
      (b) 文本形态 —— 例如"01  Plan Cover  ·  Start  →  Next: Contents"
          即"两位页码 + 章节名 + '·' + 上一页/下一页指引"的页脚行；
          页脚可能被拆成多个 run 或分布在多个文本框中，需按整框/整段文本识别。
      (c) 所有组成该页脚的 run 字号都是 7.8 磅。
      "不是 7.8 磅" → 命中扣分（任一页码处任一 run 不为 7.8 磅即命中）。
    """
    FOOTER_KEYS = ("·", "→", "Previous", "Next", "Wrap", "Start")

    def is_footer_shape(sh, prefix: str) -> bool:
        """页脚文本框判定：整框文本含页码前缀 + 至少一个页脚特征词。"""
        t = sh.text_frame.text
        if not t or prefix not in t:
            return False
        return any(k in t for k in FOOTER_KEYS)

    def collect_footer_run_sizes(sh):
        """收集页脚文本框内所有非空 run 的字号（含标签）；也覆盖'页脚被拆到多段/多 run'。"""
        entries = []  # [(size_pt_or_None, text)]
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                t = (run.text or "")
                if not t.strip():
                    continue
                sz = run.font.size.pt if run.font.size else None
                entries.append((sz, t))
        return entries

    bad = []
    for i in range(1, EXPECTED_PAGES + 1):                             # (a) 整本 PPT
        slide = prs.slides[i - 1]
        prefix = f"{i:02d}"

        # (b) 定位所有"页脚文本框"：整框文本含页码前缀 + 页脚特征词
        footer_shapes = [
            sh for sh in slide.shapes
            if sh.has_text_frame and is_footer_shape(sh, prefix)
        ]

        # 退化：单个 run 或跨 run 拼接后含"前缀 + 特征词"的文本框
        if not footer_shapes:
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                for para in sh.text_frame.paragraphs:
                    joined = "".join(r.text or "" for r in para.runs)
                    if prefix in joined and any(k in joined for k in FOOTER_KEYS):
                        footer_shapes.append(sh)
                        break

        if not footer_shapes:
            bad.append(f"第{i}页：未找到页码处文本'{prefix}  …  ·  …'")
            continue

        # (c) 收集所有组成页脚的 run 字号，逐一校验必须 = 7.8 磅
        all_entries = []
        for sh in footer_shapes:
            all_entries.extend(collect_footer_run_sizes(sh))
        if not all_entries:
            bad.append(f"第{i}页：页脚文本框存在但无可读 run 字号")
            continue

        bad_runs = [(sz, t) for sz, t in all_entries if not about(sz, 7.8)]
        if bad_runs:
            preview = "; ".join(f"'{t.strip()[:20]}'={sz}磅" for sz, t in bad_runs[:3])
            more = f"…等 {len(bad_runs)} run" if len(bad_runs) > 3 else ""
            joined_preview = "".join(t for _, t in all_entries).strip().replace("\n", " ")[:60]
            bad.append(
                f"第{i}页：页脚='{joined_preview}' 非 7.8 磅 run：{preview}{more}"
            )

    hit = len(bad) > 0
    if hit:
        show = bad[:6]
        more = f"…等共 {len(bad)} 条" if len(bad) > 6 else ""
        detail = "页码处字号不合规：\n    " + "\n    ".join(show) + ("\n    " + more if more else "")
    else:
        detail = "所有页码处字号均为 7.8 磅（覆盖多 run / 多文本框拼接的情况）"
    return hit, detail


def rule_m1_toc_font(prs):
    """-1 目录页标题如"PART 01  Project Background"字号非11.5磅，
        标题下方文本如"Policy window and regional conditions"字号非8.4磅。

    严格对齐细则的每一个点：
      (a) 目录页 —— 第 2 页；
      (b) 目录页标题 —— 形如"PART 0X  XXX"的条目标题（示例为 PART 01）；
      (c) 该标题字号是 11.5 磅；
      (d) 标题下方文本 —— 紧跟在该 PART 标题下方的说明文字（示例为
          "Policy window and regional conditions"）；
      (e) 该下方文本字号是 8.4 磅。
      标题非 11.5 磅 或 下方文本非 8.4 磅，任一不满足 → 命中扣分。
    """
    slide = prs.slides[TOC_SLIDE - 1]                                  # (a) 目录页 = 第 2 页

    # 收集该页所有非空文本框（用于按几何位置匹配"标题下方文本"）
    text_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.text_frame.text.strip():
            text_shapes.append(shape)

    def shape_sizes(sh):
        out = []
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if (run.text or "").strip():
                    out.append(run.font.size.pt if run.font.size else None)
        return out

    bad = []
    # (b) 找所有"PART 0X  XXX"标题文本框
    part_shapes = [sh for sh in text_shapes if sh.text_frame.text.strip().startswith("PART 0")]
    if not part_shapes:
        return False, "目录页未找到任何 'PART 0X' 标题，本条扣分不适用"

    for ps in part_shapes:
        title_text = ps.text_frame.text.strip().replace("\n", " ")
        title_sizes = shape_sizes(ps)
        # (c) 标题字号 = 11.5 磅（标题文本框内所有 run 字号都需为 11.5）
        title_ok = len(title_sizes) > 0 and all(about(sz, 11.5) for sz in title_sizes)

        # (d) 该 PART 标题"下方文本"：垂直方向位于其下、水平方向重叠的最近文本框
        t_top    = ps.top or 0
        t_left   = ps.left or 0
        t_right  = t_left + (ps.width or 0)
        t_bottom = t_top  + (ps.height or 0)
        below = None
        for sh in text_shapes:
            if sh is ps:
                continue
            # 跳过其它 PART 标题（避免把"PART 02"当成"PART 01"的下方文本）
            if sh.text_frame.text.strip().startswith("PART 0"):
                continue
            # 跳过页脚行
            t = sh.text_frame.text.strip()
            if "·" in t and ("→" in t or "Previous" in t or "Next" in t):
                continue
            s_top   = sh.top or 0
            s_left  = sh.left or 0
            s_right = s_left + (sh.width or 0)
            if s_top < t_bottom - (ps.height or 0) * 0.3:
                continue
            if s_right <= t_left or s_left >= t_right:
                continue
            if below is None or s_top < (below.top or 0):
                below = sh

        if below is None:
            sub_sizes = []
            sub_ok = False
            sub_text = "(未找到)"
        else:
            sub_sizes = shape_sizes(below)
            # (e) 下方文本字号 = 8.4 磅
            sub_ok = len(sub_sizes) > 0 and all(about(sz, 8.4) for sz in sub_sizes)
            sub_text = below.text_frame.text.strip().replace("\n", " ")

        if not (title_ok and sub_ok):
            bad.append(
                f"标题'{title_text[:40]}' 字号={title_sizes}(需11.5磅); "
                f"下方文本'{sub_text[:40]}' 字号={sub_sizes}(需8.4磅)"
            )

    hit = len(bad) > 0
    detail = "目录页字号不合规：\n    " + "\n    ".join(bad) if bad else "目录页所有 PART 条目字号均合规"
    return hit, detail


def rule_m5_content_pages_font(prs):
    """-5 PPT内容页（除 1、2、3、6、14、17、19、23、42、46 页，其余都为内容页）
        内容文本字号非 9.2 磅，标题文本字号非 12.5 磅加粗。

    严格对齐细则的每一个点：
      (a) 范围 —— 全 PPT 去除第 1、2、3、6、14、17、19、23、42、46 页之后的所有页（=内容页）；
      (b) 标题文本 —— 该内容页上的"标题文本框"（占位符 TITLE/CENTER_TITLE 优先，
          否则位于页顶部区域的、单行短文本框）；
      (c) 标题字号 = 12.5 磅，且加粗；
      (d) 内容文本 —— 该内容页上"正文占位符"或"主内容区文本框"；
          明确排除：页脚、页码、装饰性纯数字、极短标签（图表/坐标轴标签、编号）、
                    "Source:/Fig./Chart/图注/Note/Data:"等来源/注释类前缀；
      (e) 内容字号 = 9.2 磅；
      "标题非 12.5 磅加粗" 或 "内容非 9.2 磅"，任一不满足 → 命中扣分。
    """
    # 占位符类型枚举
    PH_TITLE, PH_CENTER_TITLE = 13, 15
    PH_SUBTITLE, PH_BODY, PH_OBJECT = 4, 2, 7

    # 装饰/注释/图表标注类前缀（大小写不敏感包含判定）
    DECOR_PREFIXES = (
        "source:", "src:", "data source", "data:", "note:", "notes:",
        "fig.", "figure ", "chart ", "table ", "ref:", "reference",
        "©", "copyright", "footnote",
    )
    FOOTER_KEYS = ("·", "→", "Previous", "Next", "Wrap")

    def ph_type(sh):
        try:
            if sh.is_placeholder and sh.placeholder_format is not None:
                pht = sh.placeholder_format.type
                return int(pht) if pht is not None else None
        except Exception:
            return None
        return None

    def is_footer_shape(sh) -> bool:
        t = sh.text_frame.text
        if not t:
            return False
        # 完整页脚形态：含 '·' + 导航词
        if "·" in t and any(k in t for k in FOOTER_KEYS):
            return True
        return False

    def is_decor_or_label(sh, page_top_half: float) -> bool:
        """判定是否为装饰/图表标注/来源说明/极短标签——不纳入正文校验。"""
        t = sh.text_frame.text.strip()
        if not t:
            return True
        # 纯数字（页码、序号、坐标轴刻度）
        stripped = t.replace("\n", " ").replace(" ", "")
        if stripped.isdigit():
            return True
        # 单个字符 / 单符号（列表标记、装饰点）
        if len(stripped) <= 1:
            return True
        # 明确的来源/图注前缀
        low = t.lower().lstrip()
        if any(low.startswith(p) for p in DECOR_PREFIXES):
            return True
        # 极短标签（<=4 字符）且位于页面上半部或角落 —— 图例/图表标签常见形态
        if len(stripped) <= 4:
            # 位于页面顶部区域（前 15%）通常是装饰性数字/角标
            if (sh.top or 0) < page_top_half * 0.3:
                return True
        return False

    def sizes_of(sh):
        out = []
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if (run.text or "").strip():
                    out.append(run.font.size.pt if run.font.size else None)
        return out

    def bolds_of(sh):
        out = []
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                if (run.text or "").strip():
                    out.append(run.font.bold)
        return out

    def max_size_of(sh):
        sz = [s for s in sizes_of(sh) if s is not None]
        return max(sz) if sz else -1

    bad = []
    for i in range(1, EXPECTED_PAGES + 1):                             # (a) 范围
        if i in NON_CONTENT_SLIDES:
            continue
        slide = prs.slides[i - 1]
        page_h = prs.slide_height or 0

        # 收集所有非空、非页脚文本框
        all_shapes = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not shape.text_frame.text.strip():
                continue
            if is_footer_shape(shape):
                continue
            all_shapes.append(shape)
        if not all_shapes:
            bad.append(f"第{i}页：无可识别文本框")
            continue

        # ---- (b) 标题定位 ----
        title_shape = None
        # 1) 占位符 TITLE / CENTER_TITLE
        for sh in all_shapes:
            if ph_type(sh) in (PH_TITLE, PH_CENTER_TITLE):
                title_shape = sh
                break
        # 2) 页面顶部（前 25%）、非装饰、非过长的单段短文本框（模板标题的典型几何）
        if title_shape is None:
            top_band = page_h * 0.25
            top_candidates = []
            for sh in all_shapes:
                if (sh.top or 0) > top_band:
                    continue
                if is_decor_or_label(sh, page_h):
                    continue
                t = sh.text_frame.text.strip()
                # 标题一般是单行、字符不超过 60
                if "\n" in t or len(t) > 60:
                    continue
                top_candidates.append(sh)
            if top_candidates:
                # 顶部候选中：优先字号最大；同字号取更靠上
                title_shape = max(
                    top_candidates,
                    key=lambda sh: (max_size_of(sh), -(sh.top or 0)),
                )
        # 3) 兜底：全页非装饰候选中字号最大且靠上
        if title_shape is None:
            non_decor = [sh for sh in all_shapes if not is_decor_or_label(sh, page_h)]
            if non_decor:
                title_shape = max(
                    non_decor,
                    key=lambda sh: (max_size_of(sh), -(sh.top or 0)),
                )

        if title_shape is None:
            bad.append(f"第{i}页：未找到标题文本框")
            continue

        title_sizes = sizes_of(title_shape)
        title_bolds = bolds_of(title_shape)
        # (c) 标题字号 = 12.5 磅 且 加粗
        title_size_ok = len(title_sizes) > 0 and all(about(sz, 12.5) for sz in title_sizes)
        title_bold_ok = len(title_bolds) > 0 and all(bool(b) for b in title_bolds)
        title_ok = title_size_ok and title_bold_ok

        # ---- (d) 正文定位 ----
        # 优先：BODY / OBJECT / SUBTITLE 占位符
        body_shapes = [
            sh for sh in all_shapes
            if sh is not title_shape and ph_type(sh) in (PH_BODY, PH_OBJECT, PH_SUBTITLE)
        ]
        # 兜底：非标题、非页脚、非装饰/图表标注/来源说明的所有文本框
        if not body_shapes:
            body_shapes = [
                sh for sh in all_shapes
                if sh is not title_shape and not is_decor_or_label(sh, page_h)
            ]

        body_sizes = []
        for sh in body_shapes:
            body_sizes.extend(sizes_of(sh))
        # (e) 内容字号 = 9.2 磅
        body_ok = len(body_sizes) > 0 and all(about(sz, 9.2) for sz in body_sizes)

        if not (title_ok and body_ok):
            body_size_set = sorted({s for s in body_sizes if s is not None})
            title_preview = title_shape.text_frame.text.strip().replace("\n", " ")[:30]
            bad.append(
                f"第{i}页：标题'{title_preview}' "
                f"字号={title_sizes} 加粗={title_bolds}(需12.5磅加粗); "
                f"正文框数={len(body_shapes)} 字号集合={body_size_set}(需9.2磅)"
            )

    hit = len(bad) > 0
    if hit:
        show = bad[:6]
        more = f"…等共 {len(bad)} 条" if len(bad) > 6 else ""
        detail = "内容页字号不合规：\n    " + "\n    ".join(show) + ("\n    " + more if more else "")
    else:
        detail = "所有内容页：标题=12.5磅加粗；内容=9.2磅（已排除页脚/图表标签/来源注释），均合规"
    return hit, detail


# ----------------------------- 主流程 -----------------------------


def _locate_document(dir_path: str):
    """在给定目录中定位被评估的 PPT 文件。
    优先使用 EXPECTED_DOC_NAME；若不存在，则取目录中任意一个 .pptx 文件。
    找不到则返回 None。
    """
    if not dir_path or not os.path.isdir(dir_path):
        return None
    expected = os.path.join(dir_path, EXPECTED_DOC_NAME)
    if os.path.isfile(expected):
        return expected
    for name in os.listdir(dir_path):
        if name.lower().endswith(".pptx") and not name.startswith("~$"):
            return os.path.join(dir_path, name)
    return None


# 加分 / 扣分规则表（label 与 rule 函数解耦，便于生成 dim2_items）
ADD_RULES = [
    (1, "第1页主标题文本框：中文“青岚县农业低空应用实施方案”翻译为“Qinglan County Agricultural Low-Altitude Application Implementation Plan”或语义一致英文。",
     "rule_p1_cover_title_translated"),
    (1, "第2页目录页文本：目录标题、PART 01至PART 07、各部分名称和说明全部为英文",
     "rule_p2_toc_all_english"),
    (3, "第3页至第46页内容文本全部为英文",
     "rule_p3_to_p46_all_english"),
]

SUB_RULES = [
    (1, "首页“Qinglan County Agricultural Low-Altitude Application Implementation Plan”字体不是32磅加粗",
     "rule_m1_cover_title_font"),
    (1, "首页“Qinglan County Agricultural Low-Altitude Application Implementation Plan”下方文本字体不是13磅",
     "rule_m1_cover_subtitle_font"),
    (5, "PPT标题页（PPT第3、6、14、17、19、23、42页）标题文本字号不是34磅，标题下方文本字号不是14磅",
     "rule_m5_title_slides_font"),
    (5, "页码处“01  Plan Cover  ·  Start  →  Next: Contents”等文本字号不是7.8磅",
     "rule_m5_page_footer_font"),
    (1, "目录页标题如“PART 01  Project Background”字号非11.5磅，标题下方文本如“Policy window and regional conditions”字号非8.4磅",
     "rule_m1_toc_font"),
    (5, "PPT内容页（除1、2、3、6、14、17、19、23、42、46页，其余都为内容页）内容文本字号非9.2磅，标题文本字号非12.5磅加粗",
     "rule_m5_content_pages_font"),
]


def evaluate(dir_path: str) -> dict:
    """统一入口：接收脚本所在目录的路径，脚本自己在该目录里定位并打开被评估文档。

    返回结构见《脚本接口差异与统一建议》§2.2。
    """
    # dim2_items 中所有条目的满分之和，用于 max_score
    max_score = sum(d for d, _, _ in ADD_RULES)

    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": max_score,
    }

    try:
        # --- 定位被评估文档 ---
        path = _locate_document(dir_path)
        if path is None:
            result["status"] = "error"
            result["error"] = f"目录 '{dir_path}' 下未找到被评估的 .pptx 文件"
            return result
        result["file_name"] = os.path.basename(path)

        # --- 维度1 ---
        ok1, reasons = check_dim1(path)
        result["dim1_pass"] = bool(ok1)
        if not ok1:
            # 未通过：给出未通过原因，dim2 不参与、总分 0
            result["dim1_reason"] = "；".join(reasons) if reasons else "维度一未通过"
            result["total_score"] = 0
            return result

        # --- 维度1 通过，进入维度2 ---
        prs = Presentation(path)
        # 规则函数在本模块中通过名字查表
        module_globals = globals()

        total = 0
        items = []

        # 加分项：命中则计入 +delta
        for delta, label, fn_name in ADD_RULES:
            fn = module_globals[fn_name]
            try:
                hit, detail = fn(prs)
            except Exception as e:  # 单条规则异常，视为未命中并在 detail 中说明
                hit, detail = False, f"规则执行异常：{e}"
            _ = detail  # detail 已计算但按需求不在结果中输出
            gained = delta if hit else 0
            total += gained
            items.append({
                "rule": label,
                "max_delta": delta,
                "delta": gained,
                "hit": bool(hit),
                "detail": "",
            })

        # 扣分项：命中则计入 -delta（max_delta 记为 -delta 以体现该项对总分的最大影响）
        for delta, label, fn_name in SUB_RULES:
            fn = module_globals[fn_name]
            try:
                hit, _detail = fn(prs)
            except Exception:
                hit = False
            lost = -delta if hit else 0
            total += lost
            items.append({
                "rule": label,
                "max_delta": -delta,
                "delta": lost,
                "hit": bool(hit),
                "detail": "",
            })

        result["dim2_items"] = items
        result["total_score"] = total
        return result

    except Exception as e:
        result["status"] = "error"
        result["error"] = f"{type(e).__name__}: {e}"
        return result


if __name__ == "__main__":
    # 本地调试：默认以脚本所在目录作为 dir_path；也可通过 argv[1] 覆盖
    target_dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(target_dir), ensure_ascii=False, indent=2))
