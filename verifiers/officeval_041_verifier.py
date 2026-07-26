# -*- coding: utf-8 -*-
"""
自动评估脚本：高速成型线碎屑分离与状态监测装置_修改完成版 (1).pptx
评估逻辑：
1. 维度1（可用与可修改性）：不满足直接判为零分
2. 维度2（完成度）：得分点+扣分点累计
"""

import os
import sys
import json

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import re
import traceback

# 脚本编号：由文件名 officeval_041_verifier.py 决定
SCRIPT_ID = "041"

# ============ 工具函数 ============
def emu_to_cm(emu):
    """EMU转厘米"""
    if emu is None:
        return None
    return emu / 914400 * 2.54

def emu_to_pt(emu):
    """EMU转磅"""
    if emu is None:
        return None
    return emu / 12700

def get_font_size_pt(run):
    """获取run的字号(磅)"""
    if run.font.size:
        return run.font.size.pt
    return None

def get_font_color_rgb(run):
    """获取run的颜色RGB字符串"""
    try:
        if run.font.color and run.font.color.rgb:
            return str(run.font.color.rgb)
    except:
        pass
    return None

def is_run_color_black(run):
    """判断 run 的字体颜色是否为黑色（放宽版）。
    同时接受：
      a) 固定 RGB：000000，或接近纯黑的极深色（R/G/B 均 <= 32）
      b) 深色主题色：tx1 / dk1 / dk2 / text1 等通常渲染为黑色的主题色
    （PowerPoint 默认正文“黑色”往往是主题色 Text 1 而非固定 000000）
    """
    return color_matches(run, '000000')


# 主题色枚举名 → 该主题色“通常渲染出的实际颜色”归类。
# 用于：当 run 的颜色是主题色(schemeClr/theme_color)而非固定 RGB 时，
# 仍能与“期望的固定 RGB 颜色”做语义匹配（放宽逻辑）。
_THEME_BLACK = ('TEXT_1', 'TEXT1', 'TX1', 'DARK_1', 'DARK1', 'DK1', 'DARK_2', 'DK2')
_THEME_WHITE = ('BACKGROUND_1', 'BACKGROUND1', 'BG1', 'LIGHT_1', 'LIGHT1', 'LT1',
                'BACKGROUND_2', 'BG2', 'LIGHT_2', 'LT2')


def _hex_to_rgb(h):
    h = str(h).upper().lstrip('#')
    if len(h) != 6:
        return None
    try:
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    except Exception:
        return None


def color_matches(run, expected_hex, tol=40):
    """放宽版颜色匹配：判断 run 字体颜色是否“算作”期望颜色 expected_hex。
    放宽规则（适用于所有页的颜色检查）：
      1) 固定 RGB：与期望色的各通道差 <= tol（默认 40）即视为同色，
         以容忍主题微调 / 取色误差 / lumMod·lumOff 造成的轻微偏差。
      2) 主题色(schemeClr / theme_color)：按“期望色的色性”归类匹配——
         期望黑 → 接受深色系主题色(tx1/dk1/dk2)；
         期望白 → 接受浅色系主题色(bg1/lt1/bg2/lt2)；
         其它彩色期望 → 主题色取不到具体 RGB 时，按色相近似(见下)从宽接受。
    """
    exp = _hex_to_rgb(expected_hex)
    if exp is None:
        return False
    er, eg, eb = exp
    try:
        color = run.font.color
    except Exception:
        return False
    if color is None:
        return False

    # 1) 固定 RGB —— 带容差比较
    try:
        rgb = color.rgb  # 主题色会抛异常
        got = _hex_to_rgb(str(rgb))
        if got is not None:
            gr, gg, gb = got
            return abs(gr - er) <= tol and abs(gg - eg) <= tol and abs(gb - eb) <= tol
    except Exception:
        pass

    # 2) 主题色 —— 按色性归类
    try:
        tc = color.theme_color
        name = str(tc).upper()
    except Exception:
        return False

    # 期望黑
    exp_is_black = (er <= 64 and eg <= 64 and eb <= 64)
    exp_is_white = (er >= 200 and eg >= 200 and eb >= 200)
    if exp_is_black:
        return any(k in name for k in _THEME_BLACK)
    if exp_is_white:
        return any(k in name for k in _THEME_WHITE)
    # 其它彩色：主题色拿不到精确 RGB，无法严格判定。
    # 从宽处理——只要是 accent 系主题色即接受（视为“设了对应主题强调色”）。
    return 'ACCENT' in name


def get_font_name(run):
    """获取字体名"""
    try:
        return run.font.name
    except:
        return None

def get_all_text_in_slide(slide):
    """获取幻灯片中所有文字"""
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    texts.append(run.text)
    if shape.has_table if hasattr(shape, 'has_table') else False:
        pass
    return texts

def get_slide_full_text(slide):
    """获取幻灯片的完整文本(包括表格)"""
    all_text = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    all_text.append(text)
        if shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    text = cell.text.strip()
                    if text:
                        all_text.append(text)
    return all_text

def get_slide_text_joined(slide):
    """获取幻灯片所有文本拼接"""
    return " ".join(get_slide_full_text(slide))

def shape_has_text(shape, target):
    """检查shape是否包含目标文字"""
    if shape.has_text_frame:
        full = shape.text_frame.text
        if target in full:
            return True
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if target in cell.text:
                    return True
    return False

def slide_contains_text(slide, target):
    """检查幻灯片是否包含目标文字(忽略换行)"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            # 把换行替换掉再匹配
            full = shape.text_frame.text.replace('\n', '').replace('\r', '').replace(' ', '')
            if target.replace(' ', '') in full:
                return True
            # 也按原始方式检查
            if target in shape.text_frame.text:
                return True
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if target in cell.text:
                        return True
    return False

def count_images_in_slide(slide):
    """统计幻灯片中图片数量"""
    count = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            count += 1
        # 也检查placeholder中的图片
        try:
            if hasattr(shape, 'image'):
                count += 1
        except:
            pass
    # 去重：shape_type==13的已经有image属性
    # 重新统计
    count = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:
            count += 1
    return count

def get_images_in_slide(slide):
    """获取幻灯片中所有图片shape"""
    imgs = []
    for shape in slide.shapes:
        if shape.shape_type == 13:
            imgs.append(shape)
    return imgs

def has_video_in_slide(slide):
    """检查幻灯片是否有视频对象"""
    for shape in slide.shapes:
        # 检查是否是媒体占位符或嵌入视频
        if shape.shape_type == 3:  # MEDIA
            return True
        # 通过XML检查video
        sp_xml = shape._element.xml
        if 'video' in sp_xml.lower() or 'media' in sp_xml.lower():
            return True
    # 通过slide的relationship检查
    try:
        for rel in slide.part.rels.values():
            if 'video' in rel.reltype.lower() or 'media' in rel.reltype.lower():
                return True
    except:
        pass
    return False

def check_video_autoplay(slide):
    """检查视频是否设为自动播放"""
    slide_xml = slide._element.xml
    tree = etree.fromstring(slide_xml.encode('utf-8') if isinstance(slide_xml, str) else slide_xml)
    timing = tree.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}timing')
    if timing is None:
        return False
    xml_str = etree.tostring(timing, encoding='unicode')
    has_main_seq = 'nodeType="mainSeq"' in xml_str
    if not has_main_seq:
        return False
    # 方式1: afterPrevious/withPrevious 触发
    if 'afterPrevious' in xml_str or 'withPrevious' in xml_str:
        return True
    # 方式2: delay="0" 且含 mediacall/playFrom (进页自动播放)
    if 'delay="0"' in xml_str and ('mediacall' in xml_str or 'playFrom' in xml_str):
        return True
    return False

def check_animation_effect(slide, effect_name):
    """检查幻灯片是否有指定动画效果，不能把任意动画构建列表当作指定效果。"""
    slide_xml = slide._element.xml
    # 动画效果名映射
    effect_map = {
        '百叶窗': ['blinds', 'blind'],
        '轮辐': ['wheel', 'spoke'],
        '涟漪': ['ripple'],
        '溶解': ['dissolve'],
    }
    targets = effect_map.get(effect_name, [effect_name.lower()])
    xml_lower = slide_xml.lower()
    for t in targets:
        if t in xml_lower:
            return True
    # 涟漪在 PowerPoint 中常以"动作路径"预设实现(presetClass="path" presetID="24")，
    # XML 里不会出现字面量 ripple，需额外按动画节点结构识别。
    if effect_name == '涟漪':
        if 'presetid="24"' in xml_lower and 'presetclass="path"' in xml_lower:
            return True
    return False

def check_animation_on_image(slide, effect_name):
    """检查动画是否作用在图片对象上"""
    # 简化：只要有动画且有图片就认为满足
    has_anim = check_animation_effect(slide, effect_name)
    has_img = len(get_images_in_slide(slide)) > 0
    return has_anim and has_img

def get_table_shapes(slide):
    """获取幻灯片中的表格shape"""
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            tables.append(shape)
    return tables

def get_text_shapes_with_info(slide):
    """获取幻灯片中所有文本框的详细信息"""
    results = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    info = {
                        'text': run.text,
                        'font_name': get_font_name(run),
                        'font_size_pt': get_font_size_pt(run),
                        'font_color': get_font_color_rgb(run),
                        'shape_left_cm': emu_to_cm(shape.left),
                        'shape_top_cm': emu_to_cm(shape.top),
                        'shape_width_cm': emu_to_cm(shape.width),
                        'shape_height_cm': emu_to_cm(shape.height),
                        'alignment': para.alignment,
                        'bold': run.font.bold,
                    }
                    results.append(info)
    return results

def find_text_runs_in_slide(slide, target_text):
    """在幻灯片中找到包含目标文字的run信息"""
    results = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                para_text = para.text
                if target_text in para_text:
                    for run in para.runs:
                        if target_text in run.text:
                            results.append({
                                'run': run,
                                'para': para,
                                'shape': shape,
                                'font_size_pt': get_font_size_pt(run),
                                'font_color': get_font_color_rgb(run),
                                'font_name': get_font_name(run),
                            })
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        if target_text in para.text:
                            for run in para.runs:
                                if target_text in run.text:
                                    results.append({
                                        'run': run,
                                        'para': para,
                                        'shape': shape,
                                        'font_size_pt': get_font_size_pt(run),
                                        'font_color': get_font_color_rgb(run),
                                        'font_name': get_font_name(run),
                                    })
    return results

def check_shape_overlap(shape1, shape2):
    """检查两个shape是否重叠"""
    l1, t1, r1, b1 = shape1.left, shape1.top, shape1.left + shape1.width, shape1.top + shape1.height
    l2, t2, r2, b2 = shape2.left, shape2.top, shape2.left + shape2.width, shape2.top + shape2.height
    return not (r1 <= l2 or r2 <= l1 or b1 <= t2 or b2 <= t1)

def get_slide_width_height(prs):
    """获取幻灯片宽高(cm)"""
    w = emu_to_cm(prs.slide_width)
    h = emu_to_cm(prs.slide_height)
    return w, h

def check_table_text_alignment(slide, expected_align=PP_ALIGN.LEFT):
    """检查表格文字对齐方式, None视为左对齐(默认)"""
    for shape in slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        if para.text.strip():
                            # None = 继承 = 通常是左对齐
                            if para.alignment is not None and para.alignment != expected_align:
                                return False
    return True

def check_image_covers_text(slide, texts_to_check):
    """检查图片是否遮挡指定文字"""
    imgs = get_images_in_slide(slide)
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for t in texts_to_check:
                if t in shape.text_frame.text:
                    text_shapes.append(shape)
                    break
    for img in imgs:
        for ts in text_shapes:
            if check_shape_overlap(img, ts):
                return True
    return False

def check_textbox_has_yellow_bg(slide):
    """检查是否有文本框有多余黄色背景"""
    for shape in slide.shapes:
        if shape.has_text_frame:
            try:
                xml = shape._element.xml
                # 检查shape本身有solidFill为黄色
                # 排除圆形等图形，只检查文本框(sp:txBody的父级)
                if '<a:solidFill>' in xml:
                    import re
                    fills = re.findall(r'<a:solidFill>\s*<a:srgbClr val="([^"]+)"', xml)
                    for f in fills:
                        if f.upper() in ['FFFF00', 'FFD700', 'FFC000', 'FACC15', 'FEF08A', 'FBBF24']:
                            # 确认这是文本框的背景而非文字颜色
                            # 通过检查是否在spPr中(shape属性=背景填充)
                            sp_pr_match = re.search(r'<p:spPr[^>]*>(.*?)</p:spPr>', xml, re.DOTALL)
                            if sp_pr_match and f in sp_pr_match.group(1):
                                return True
            except:
                pass
    return False

def is_image_type_ppt(slide):
    """检查页面是否是图片型(不可编辑)"""
    imgs = get_images_in_slide(slide)
    if not imgs:
        return False
    # 如果一张图片占据了90%以上面积
    for img in imgs:
        if img.width and img.height:
            slide_w = slide.slide_layout.slide_master.slide_width if hasattr(slide, 'slide_layout') else None
            # 简化判断
            pass
    return False

def check_large_image_blocking(slide, prs):
    """检查是否有超大图片(>90%页面)且标题不可编辑"""
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    threshold_w = slide_w * 0.9
    threshold_h = slide_h * 0.9
    imgs = get_images_in_slide(slide)
    for img in imgs:
        if img.width and img.height:
            if img.width > threshold_w and img.height > threshold_h:
                # 检查是否有可编辑的中文标题
                has_editable_title = False
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != img:
                        text = shape.text_frame.text
                        if len(text) > 1 and any('一' <= c <= '鿿' for c in text):
                            has_editable_title = True
                            break
                if not has_editable_title:
                    return True
    return False

def check_content_overflow(slide, prs):
    """检查是否有内容超出页面边界"""
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    for shape in slide.shapes:
        if shape.left is not None and shape.width is not None:
            if shape.left + shape.width > slide_w * 1.02:  # 2%容差
                return True
        if shape.top is not None and shape.height is not None:
            if shape.top + shape.height > slide_h * 1.02:
                return True
    return False


# ============ 主评估逻辑 ============
def evaluate(dir_path: str) -> dict:
    """在给定目录中定位被评估的 PPT 文件并执行评估,返回结构化字典。"""
    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": True,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": 0,
    }

    # 在给定目录中定位被评估的 .pptx 文件（只识别 .pptx，不再兼容老版 .ppt 二进制格式）
    pptx_file = None
    try:
        if os.path.isdir(dir_path):
            for name in os.listdir(dir_path):
                if name.startswith('~$'):
                    continue
                if name.lower().endswith('.pptx'):
                    pptx_file = os.path.join(dir_path, name)
                    break
    except Exception as e:
        result["status"] = "error"
        result["error"] = f"读取目录失败: {e}"
        return result
    if pptx_file is None:
        result["status"] = "error"
        result["error"] = f"目录中未找到 .pptx 文件: {dir_path}"
        return result
    result["file_name"] = os.path.basename(pptx_file)

    # ========== 维度1：可用与可修改性 ==========
    dim1_pass = True
    dim1_reasons = []

    # 1.1 文件格式检查
    if not pptx_file.lower().endswith('.pptx'):
        dim1_pass = False
        dim1_reasons.append("文件不是.pptx格式")

    # 1.2 文件可正常打开
    try:
        prs = Presentation(pptx_file)
    except Exception as e:
        result["dim1_pass"] = False
        result["dim1_reason"] = f"文件无法正常打开: {e}"
        return result

    # 1.3 页数检查
    slide_count = len(prs.slides)
    if slide_count != 23:
        dim1_pass = False
        dim1_reasons.append(f"页数为{slide_count}页，要求23页")

    if not dim1_pass:
        result["dim1_pass"] = False
        result["dim1_reason"] = "; ".join(dim1_reasons)
        return result

    # ========== 维度2：完成度评分 ==========
    scores = []  # (分值, 描述, 是否满足)
    slides = list(prs.slides)

    # ---- +5: 页面顺序调整 ----
    # 细则要求：第1页~第23页依次为指定标题
    # 第1页：封面"高速成型线碎屑分离与状态监测装置改进汇报"
    # 第2页：目录 CONTENTS
    # 第3~23页：依次为下列21个标题
    expected_titles = [
        ["高速成型线碎屑分离与状态监测装置改进汇报"],                # 第1页
        ["目录", "CONTENTS"],                                          # 第2页(需同时含"目录"与"CONTENTS")
        ["项目定位"],                                                  # 第3页
        ["风险机制", "三个要素相互叠加"],                              # 第4页
        ["改造前", "回收路径长", "观察节点少"],                        # 第5页
        ["改造目标", "把风险链拆开"],                                  # 第6页
        ["需求转换", "从现场声音到关键需求"],                          # 第7页
        ["创新借鉴", "分割与预先感知"],                                # 第8页
        ["目标设定", "用可量化指标约束设计"],                          # 第9页
        ["可行性论证", "从工艺", "技术", "资源三方面评估"],            # 第10页
        ["模拟试验一", "筛分路径验证"],                                # 第11页
        ["模拟试验二", "粉料含量与异常蔓延趋势"],                      # 第12页
        ["总体方案", "物理分路", "在线状态感知"],                      # 第13页
        ["研发指标", "功能", "性能与兼容性"],                          # 第14页
        ["制定对策", "任务", "目标与交付物"],                          # 第15页
        ["对策实施一", "筛分系统设计与加工"],                          # 第16页
        ["对策实施二", "可开合观察窗制作"],                            # 第17页
        ["对策实施三", "监测画面接入操作位"],                          # 第18页
        ["整体试运行", "连续班次结果"],                                # 第19页
        ["效果验证", "风险链条被削弱"],                                # 第20页
        ["标准化与巩固措施"],                                          # 第21页
        ["总结与展望"],                                                # 第22页
        ["感谢聆听"],                                                  # 第23页
    ]
    order_ok = True
    if len(slides) != len(expected_titles):
        order_ok = False
    else:
        for i, slide in enumerate(slides):
            text_norm = get_slide_text_joined(slide).replace(" ", "").replace("\n", "")
            for kw in expected_titles[i]:
                if kw.replace(" ", "") not in text_norm:
                    order_ok = False
                    break
            if not order_ok:
                break
    scores.append((5, "页面顺序调整正确", order_ok))

    # ---- +5: 章节标识文字组 ----
    # 细则要求：
    #   第3、4、5页顶部右侧 → "01 背景"
    #   第6页顶部右侧 → "02 目标"
    #   第7、8页顶部右侧 → "03 需求"
    #   第9、10页顶部右侧 → "04 目标"
    #   第11、12页顶部右侧 → "05 验证"
    #   第13、14页顶部右侧 → "06 方案"
    #   第15、16、17、18页顶部右侧 → "07 实施"
    #   第19、20页顶部右侧 → "08 验证"
    #   第21页顶部右侧 → "09 巩固"
    # 位置：距左 27cm–32.5cm、距上 0.8cm–1.8cm
    # 字体：微软雅黑;字号：20磅;颜色：00A6A6 (R=0 G=166 B=166)
    chapter_map = {
        # key 为页码(1-indexed),对应 slides 索引为 key-1
        3: "01 背景", 4: "01 背景", 5: "01 背景",
        6: "02 目标",
        7: "03 需求", 8: "03 需求",
        9: "04 目标", 10: "04 目标",
        11: "05 验证", 12: "05 验证",
        13: "06 方案", 14: "06 方案",
        15: "07 实施", 16: "07 实施", 17: "07 实施", 18: "07 实施",
        19: "08 验证", 20: "08 验证",
        21: "09 巩固",
    }
    chapter_all_ok = True
    for page_no, expected_text in chapter_map.items():
        slide_idx = page_no - 1
        if slide_idx >= len(slides):
            chapter_all_ok = False
            break
        slide = slides[slide_idx]
        expected_norm = expected_text.replace(" ", "")
        page_found = False
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run_norm = run.text.replace(" ", "")
                    if expected_norm not in run_norm:
                        continue
                    page_found = True
                    # 位置：距左 27~32.5cm、距上 0.8~1.8cm
                    left_cm = emu_to_cm(shape.left)
                    top_cm = emu_to_cm(shape.top)
                    if left_cm is None or not (27 <= left_cm <= 32.5):
                        chapter_all_ok = False
                    if top_cm is None or not (0.8 <= top_cm <= 1.8):
                        chapter_all_ok = False
                    # 字体：微软雅黑
                    fname = get_font_name(run)
                    if fname is None or "微软雅黑" not in fname:
                        chapter_all_ok = False
                    # 字号：20磅
                    sz = get_font_size_pt(run)
                    if sz is None or abs(sz - 20) > 0.5:
                        chapter_all_ok = False
                    # 颜色：00A6A6
                    if not color_matches(run, '00A6A6'):
                        chapter_all_ok = False
                    break
                if page_found:
                    break
            if page_found:
                break
        if not page_found:
            chapter_all_ok = False
    scores.append((5, "章节标识文字组正确(位置/字体/颜色)", chapter_all_ok))

    # ---- +1: 第2页目录页标题文字 ----
    # 细则要求：
    #   1) 第2页左上或中上位置出现"目录"和"CONTENTS"
    #   2) "目录" — 微软雅黑、40磅、白色 FFFFFF
    #   3) "CONTENTS" — 英文无衬线字体、约10磅、加粗、浅绿色 99F6E4
    #   4) 目录标题不遮挡六个目录板块
    s2 = slides[1]
    s2_text = get_slide_text_joined(s2)
    has_mulu = "目录" in s2_text
    has_contents = "CONTENTS" in s2_text.upper()
    s2_title_ok = has_mulu and has_contents

    slide_w_emu = prs.slide_width
    slide_h_emu = prs.slide_height

    def _shape_of_text(slide, target):
        for shape in slide.shapes:
            if shape.has_text_frame and target in shape.text_frame.text:
                return shape
        return None

    # 位置：左上或中上 → top 在页面上半部(<=50%);left 在 0 ~ 75% 之间
    mulu_shape = _shape_of_text(s2, "目录")
    contents_shape = None
    for shape in s2.shapes:
        if shape.has_text_frame and "CONTENTS" in shape.text_frame.text.upper():
            contents_shape = shape
            break

    if s2_title_ok:
        for sh in (mulu_shape, contents_shape):
            if sh is None or sh.top is None or sh.left is None:
                s2_title_ok = False
                break
            if sh.top > slide_h_emu * 0.5:
                s2_title_ok = False
                break
            if sh.left < 0 or sh.left > slide_w_emu * 0.75:
                s2_title_ok = False
                break

    # "目录" 字体微软雅黑、40磅、白色 FFFFFF
    if s2_title_ok:
        runs = find_text_runs_in_slide(s2, "目录")
        if not runs:
            s2_title_ok = False
        for r in runs:
            fname = r['font_name']
            if fname is None or "微软雅黑" not in fname:
                s2_title_ok = False
            if r['font_size_pt'] is None or abs(r['font_size_pt'] - 40) > 0.5:
                s2_title_ok = False
            if r['font_color'] is None or r['font_color'].upper() != 'FFFFFF':
                s2_title_ok = False
            break

    # "CONTENTS" 英文无衬线字体、约10磅、加粗、浅绿色 99F6E4
    sans_serif_fonts = {
        'arial', 'calibri', 'helvetica', 'verdana', 'tahoma', 'segoe ui',
        'trebuchet ms', 'century gothic', 'microsoft yahei', 'microsoft yahei ui',
        '微软雅黑', 'dejavu sans', 'liberation sans', 'noto sans', 'roboto',
        'open sans', 'source sans pro', 'pingfang sc', 'pingfang', '苹方',
    }
    if s2_title_ok:
        runs = find_text_runs_in_slide(s2, "CONTENTS")
        if not runs:
            s2_title_ok = False
        for r in runs:
            fname = r['font_name']
            if fname is None:
                s2_title_ok = False
            else:
                fname_norm = fname.strip().lower()
                if not any(s in fname_norm for s in sans_serif_fonts):
                    s2_title_ok = False
            if r['font_size_pt'] is None or abs(r['font_size_pt'] - 10) > 1.5:
                s2_title_ok = False
            # 加粗
            try:
                if not bool(r['run'].font.bold):
                    s2_title_ok = False
            except Exception:
                s2_title_ok = False
            if r['font_color'] is None or r['font_color'].upper() != '99F6E4':
                s2_title_ok = False
            break

    # 标题不遮挡六个目录板块
    s2_section_keywords = ["项目背景", "目标设定", "方案设计", "实施过程", "效果验证", "标准化"]
    if s2_title_ok:
        section_shapes = []
        for shape in s2.shapes:
            if not shape.has_text_frame:
                continue
            for kw in s2_section_keywords:
                if kw in shape.text_frame.text:
                    section_shapes.append(shape)
                    break
        for title_shape in (mulu_shape, contents_shape):
            if title_shape is None:
                continue
            for sec_shape in section_shapes:
                if title_shape is sec_shape:
                    continue
                try:
                    if check_shape_overlap(title_shape, sec_shape):
                        s2_title_ok = False
                        break
                except Exception:
                    pass
            if not s2_title_ok:
                break
    scores.append((1, "第2页目录页标题文字", s2_title_ok))

    # ---- +3: 第2页目录六个板块文字 ----
    # 细则要求：
    #   1) 出现六个板块文字（"01 项目背景与需求" 等，编号与标题可分别位于不同 shape）
    #   2) 六个板块标题字体为微软雅黑、30磅、白色 FFFFFF
    #   3) 保持纵向或分组排列结构
    #   4) 文字之间无明显重叠
    # 说明：编号"01"~"06"与标题部分常作为两个独立文本框分列设计——这里允许
    #       编号与标题分布在两个 shape 中,只要：
    #         a) 整页文本里两段文字都能找到；
    #         b) 编号 shape 与标题 shape 在视觉上邻近(同一行,水平相邻)；
    #         c) 字体/字号/颜色校验作用在"标题部分"的 run 上(标题才是细则约束的对象)。
    s2_sections = [("01", "项目背景与需求"),
                   ("02", "目标设定与可行性"),
                   ("03", "方案设计与优选"),
                   ("04", "实施过程"),
                   ("05", "效果验证"),
                   ("06", "标准化与展望")]
    s2_sections_ok = True
    section_shapes_for_layout = []  # 用于布局/重叠校验的"标题"shape 列表

    def _find_shape_containing(slide, text):
        """在 slide 中找到 text_frame 含 target 的 shape (忽略空格/换行)"""
        target_norm = text.replace(" ", "").replace("\n", "")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf_norm = shape.text_frame.text.replace(" ", "").replace("\n", "")
            if target_norm in tf_norm:
                return shape
        return None

    for num, title in s2_sections:
        # (1) 整页文本里两段都存在;允许编号和标题在同一 shape 或分两个 shape
        combined_shape = _find_shape_containing(s2, num + title)  # 合并在一起的情况
        num_shape = _find_shape_containing(s2, num) if combined_shape is None else combined_shape
        title_shape = _find_shape_containing(s2, title) if combined_shape is None else combined_shape

        if title_shape is None or num_shape is None:
            s2_sections_ok = False
            break

        # 若编号与标题分两个 shape,要求二者视觉上同一行(垂直中心差距 < 标题 shape 高度)
        if num_shape is not title_shape:
            try:
                num_cy = num_shape.top + num_shape.height / 2
                title_cy = title_shape.top + title_shape.height / 2
                if abs(num_cy - title_cy) > max(num_shape.height, title_shape.height):
                    s2_sections_ok = False
                    break
            except Exception:
                s2_sections_ok = False
                break

        section_shapes_for_layout.append(title_shape)

        # (2) 字体/字号/颜色 — 作用在标题部分的 run 上
        title_norm = title.replace(" ", "")
        title_run = None
        for para in title_shape.text_frame.paragraphs:
            if title_norm not in para.text.replace(" ", ""):
                continue
            for run in para.runs:
                if not run.text.strip():
                    continue
                # 优先取包含标题首字的 run
                if title[0] in run.text or run.text.strip() in title:
                    title_run = run
                    break
            if title_run is None:
                # 兜底：取段落首个非空 run
                for run in para.runs:
                    if run.text.strip():
                        title_run = run
                        break
            if title_run:
                break

        if title_run is None:
            s2_sections_ok = False
            break

        fname = get_font_name(title_run)
        if fname is None or "微软雅黑" not in fname:
            s2_sections_ok = False
            break
        sz = get_font_size_pt(title_run)
        if sz is None or abs(sz - 30) > 0.5:
            s2_sections_ok = False
            break
        if not color_matches(title_run, 'FFFFFF'):
            s2_sections_ok = False
            break

    # (3) 排列结构：纵向（top 严格递增）或分组（如 2 列 × 3 行 / 3 列 × 2 行）
    if s2_sections_ok and len(section_shapes_for_layout) == 6:
        tops = [sh.top for sh in section_shapes_for_layout]
        lefts = [sh.left for sh in section_shapes_for_layout]
        tops_sorted = sorted(tops)
        is_vertical = tops_sorted == tops and (max(tops) - min(tops)) > 0
        unique_lefts = len(set(round(l / 100000) for l in lefts))
        unique_tops = len(set(round(t / 100000) for t in tops))
        is_grouped = unique_lefts >= 2 and unique_tops >= 2
        if not (is_vertical or is_grouped):
            s2_sections_ok = False

    # (4) 文字之间无明显重叠（仅校验六个标题 shape 之间）
    if s2_sections_ok and len(section_shapes_for_layout) == 6:
        # 去重：编号和标题合并在同一个 shape 的情况下不重复比较
        unique_shapes = []
        seen_ids = set()
        for sh in section_shapes_for_layout:
            if id(sh) in seen_ids:
                continue
            seen_ids.add(id(sh))
            unique_shapes.append(sh)
        for i in range(len(unique_shapes)):
            for j in range(i + 1, len(unique_shapes)):
                try:
                    if check_shape_overlap(unique_shapes[i], unique_shapes[j]):
                        s2_sections_ok = False
                        break
                except Exception:
                    pass
            if not s2_sections_ok:
                break
    scores.append((3, "第2页目录六个板块文字", s2_sections_ok))

    # ---- +3: 第4页三张图片插入 ----
    # 细则要求（已按用户建议删除"内容依次为热成像地面热点图、白色粉状物料堆、金属
    #  设备内碎屑分离场景"一条 —— 该条属于图片语义识别,在不引入 COM/图像识别的
    #  前提下无法自动判定；本脚本坚持不使用 COM,故不校验图片内容）：
    #   1) 第4页出现三张图片
    #   2) 三张图片按从左到右或从上到下顺序排列
    #   3) 单张图片横纵比约 4:3 (≈1.333)
    #   4) 均带淡蓝色单实线边框,线宽 0.75–1.5 磅
    #   5) 不遮挡"粉状碎屑""残留热斑""封闭堆积"等页面原有文字
    s4 = slides[3]
    s4_imgs = get_images_in_slide(s4)
    s4_img_ok = len(s4_imgs) >= 3
    imgs3 = s4_imgs[:3]

    # (2) 排列顺序：从左到右 或 从上到下
    if s4_img_ok:
        lefts = [img.left for img in imgs3]
        tops = [img.top for img in imgs3]
        left_sorted = all(lefts[i] < lefts[i + 1] for i in range(len(lefts) - 1))
        top_sorted = all(tops[i] < tops[i + 1] for i in range(len(tops) - 1))
        if not (left_sorted or top_sorted):
            s4_img_ok = False

    # (3) 横纵比约 4:3 (容差 ±15%,即 1.13 ~ 1.53)
    if s4_img_ok:
        for img in imgs3:
            w = img.width
            h = img.height
            if not (w and h and h > 0):
                s4_img_ok = False
                break
            ratio = w / h
            if not (4 / 3 * 0.85 <= ratio <= 4 / 3 * 1.15):
                s4_img_ok = False
                break

    # (4) 边框:淡蓝色单实线 + 线宽 0.75–1.5 磅 (1磅=12700EMU → 9525 ~ 19050 EMU)
    light_blue_hex = {
        'ADD8E6', 'B0E0E6', 'BFDBFE', '93C5FD', '60A5FA',
        'A5C8E1', 'BDD7EE', 'D9E1F2', 'CFE2F3', '9DC3E6',
        'A6CEE3', '7FB3D5', '85C1E9', 'AED6F1', 'D6EAF8',
        'C6DBEF', 'DEEBF7', 'B4C7E7', '8FAADC', '2E75B6',
    }
    def _is_light_blue(hex_color):
        if hex_color is None:
            return False
        h = hex_color.upper()
        if h in light_blue_hex:
            return True
        try:
            r = int(h[0:2], 16); g = int(h[2:4], 16); b = int(h[4:6], 16)
        except Exception:
            return False
        # 放宽判定：只要属于“蓝色系且偏亮(淡)”的任意一种淡蓝都接受。
        #   1) 蓝色分量主导：b 不小于 r、g（允许青蓝/天蓝，g 略高于 b 也可）
        #   2) 整体偏亮(淡)：蓝分量足够高，或整体亮度足够
        blue_dominant = (b >= r - 10) and (b >= g - 40)
        light = (b >= 120) and ((r + g + b) >= 280)
        return blue_dominant and light

    if s4_img_ok:
        for img in imgs3:
            xml = img._element.xml
            ln_match = re.search(r'<a:ln\b([^>]*)>(.*?)</a:ln>', xml, re.DOTALL)
            if not ln_match:
                s4_img_ok = False
                break
            ln_attrs = ln_match.group(1)
            ln_body = ln_match.group(2)
            # 线宽
            w_attr = re.search(r'w="(\d+)"', ln_attrs)
            if not w_attr:
                s4_img_ok = False
                break
            ln_w = int(w_attr.group(1))
            if not (9525 <= ln_w <= 19050):
                s4_img_ok = False
                break
            # 单实线: <a:solidFill> 存在,且无 <a:prstDash> 或为 solid
            if '<a:solidFill>' not in ln_body:
                s4_img_ok = False
                break
            dash_match = re.search(r'<a:prstDash[^>]*val="([^"]+)"', ln_body)
            if dash_match and dash_match.group(1).lower() != 'solid':
                s4_img_ok = False
                break
            # 颜色：淡蓝色
            # 同时接受两种表达方式：
            #   a) 固定 RGB：<a:srgbClr val="ADD8E6"/> 等淡蓝色值
            #   b) 主题色：<a:schemeClr val="accent1"> + lumMod/lumOff 变浅，
            #      只要色相属于蓝色系主题色即视为淡蓝（PowerPoint“浅色变体”常见做法）
            fill_match = re.search(r'<a:solidFill>(.*?)</a:solidFill>', ln_body, re.DOTALL)
            fill_body = fill_match.group(1) if fill_match else ''
            color_ok = False
            srgb = re.search(r'<a:srgbClr val="([^"]+)"', fill_body)
            if srgb and _is_light_blue(srgb.group(1)):
                color_ok = True
            scheme = re.search(r'<a:schemeClr val="([^"]+)"', fill_body)
            if scheme:
                scheme_val = scheme.group(1).lower()
                # 蓝色系主题色（不同模板里蓝色多落在 accent1/accent5/tx2/dk2）
                blue_scheme = {'accent1', 'accent5', 'tx2', 'dk2'}
                has_light = ('lummod' in fill_body.lower()) or ('lumoff' in fill_body.lower()) \
                    or ('tint' in fill_body.lower())
                # 蓝色系主题色 + 变浅处理 → 视为淡蓝；
                # 即便没有显式变浅，蓝色系主题色也按淡蓝接受（评分从宽）
                if scheme_val in blue_scheme:
                    color_ok = True
            if not color_ok:
                s4_img_ok = False
                break

    # (5) 不遮挡原有文字
    if s4_img_ok:
        if check_image_covers_text(s4, ["粉状碎屑", "残留热斑", "封闭堆积"]):
            s4_img_ok = False
    scores.append((3, "第4页三张图片插入(排列/比例/边框/不遮挡)", s4_img_ok))

    # ---- +1: 第5页表格与视频设置 ----
    # ---- +1: 第5页表格与视频设置 ----
    # 细则要求：
    #   1) 页面中部表格
    #   2) 表格内文字字体为微软雅黑
    #   3) 字号 10磅
    #   4) 表格文字仍可编辑（即真实表格 shape,而非图片型）
    s5 = slides[4]
    s5_table_ok = False
    slide_h_emu_s5 = prs.slide_height
    for shape in s5.shapes:
        if not shape.has_table:
            continue
        # (1) 中部：表格垂直中心位于页面 20%~80% 之间
        if shape.top is None or shape.height is None:
            continue
        center = shape.top + shape.height / 2
        if not (slide_h_emu_s5 * 0.2 <= center <= slide_h_emu_s5 * 0.8):
            continue
        # (4) 真实表格 shape — has_table 为 True 即可编辑表格;遍历其 run 校验字体/字号
        table_runs_ok = True
        has_any_run = False
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_any_run = True
                        # (2) 字体微软雅黑
                        fname = get_font_name(run)
                        if fname is None or "微软雅黑" not in fname:
                            table_runs_ok = False
                        # (3) 字号 10 磅
                        sz = get_font_size_pt(run)
                        if sz is None or abs(sz - 10) > 0.5:
                            table_runs_ok = False
        if has_any_run and table_runs_ok:
            s5_table_ok = True
            break
    scores.append((1, "第5页表格文字为微软雅黑10磅可编辑", s5_table_ok))

    # ---- +5: 第5页视频自动播放 ----
    # 细则要求：
    #   1) 第5页左侧页面中的视频对象
    #   2) 设置为进入本页后自动播放
    #   3) 视频对象仍位于页面主体区域内
    #   4) 没有遮挡标题和表格
    slide_w_emu_s5 = prs.slide_width
    slide_h_emu_s5b = prs.slide_height

    def _find_video_shapes(slide):
        vids = []
        for shape in slide.shapes:
            try:
                if shape.shape_type == 3:  # MEDIA
                    vids.append(shape)
                    continue
                sp_xml = shape._element.xml
                if 'video' in sp_xml.lower() or '<p:videoFile' in sp_xml or 'media' in sp_xml.lower():
                    vids.append(shape)
            except Exception:
                pass
        return vids

    s5_video_shapes = _find_video_shapes(s5)
    s5_video_ok = len(s5_video_shapes) > 0 and check_video_autoplay(s5)

    # (1) 视频对象位于页面左侧：视频中心 left+width/2 在页面宽度 0%~50% 区间
    if s5_video_ok:
        on_left = False
        for v in s5_video_shapes:
            if v.left is None or v.width is None:
                continue
            cx = v.left + v.width / 2
            if 0 <= cx <= slide_w_emu_s5 * 0.5:
                on_left = True
                break
        if not on_left:
            s5_video_ok = False

    # (3) 位于页面主体区域内：视频矩形不超出页面边界
    if s5_video_ok:
        in_bounds = False
        for v in s5_video_shapes:
            if v.left is None or v.top is None or v.width is None or v.height is None:
                continue
            if (v.left >= 0 and v.top >= 0
                    and v.left + v.width <= slide_w_emu_s5
                    and v.top + v.height <= slide_h_emu_s5b):
                in_bounds = True
                break
        if not in_bounds:
            s5_video_ok = False

    # (4) 不遮挡标题和表格
    if s5_video_ok:
        # 标题：第5页应为"改造前：回收路径长、观察节点少"或顶部标题 shape
        title_shapes = []
        for shape in s5.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text
            if "改造前" in txt or "回收路径" in txt or "观察节点" in txt:
                title_shapes.append(shape)
        table_shapes = [sh for sh in s5.shapes if sh.has_table]
        cover_bad = False
        for v in s5_video_shapes:
            for ts in title_shapes + table_shapes:
                try:
                    if check_shape_overlap(v, ts):
                        cover_bad = True
                        break
                except Exception:
                    pass
            if cover_bad:
                break
        if cover_bad:
            s5_video_ok = False
    scores.append((5, "第5页视频自动播放", s5_video_ok))

    # ---- +1: 第6页目标关键词 ----
    # 细则要求：
    #   1) "分离""观察""维护"三个词均出现
    #   2) 字体均为微软雅黑
    #   3) 字号均为15磅
    #   4) 三个词分别位于三个模块的主要标题位置（三个不同 shape 中）
    #   5) 颜色分别为：分离=00A6A6,观察=1D4ED8,维护=F97316
    s6 = slides[5]
    s6_keywords = ["分离", "观察", "维护"]
    s6_colors = ["00A6A6", "1D4ED8", "F97316"]
    s6_kw_ok = True
    s6_keyword_shapes = []
    for kw, expected_color in zip(s6_keywords, s6_colors):
        found_run = None
        found_shape = None
        for shape in s6.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if kw in run.text:
                        found_run = run
                        found_shape = shape
                        break
                if found_run:
                    break
            if found_run:
                break
        if found_run is None:
            s6_kw_ok = False
            break
        s6_keyword_shapes.append(found_shape)
        # (2) 字体：微软雅黑
        fname = get_font_name(found_run)
        if fname is None or "微软雅黑" not in fname:
            s6_kw_ok = False
            break
        # (3) 字号：15磅
        sz = get_font_size_pt(found_run)
        if sz is None or abs(sz - 15) > 0.5:
            s6_kw_ok = False
            break
        # (5) 颜色
        if not color_matches(found_run, expected_color):
            s6_kw_ok = False
            break

    # (4) 三个词分别位于三个不同的 shape（即三个模块）
    if s6_kw_ok:
        if len(s6_keyword_shapes) != 3 or len(set(id(sh) for sh in s6_keyword_shapes)) != 3:
            s6_kw_ok = False
    scores.append((1, "第6页目标关键词文字", s6_kw_ok))

    # ---- +1: 第7页需求表格文字 ----
    # 细则要求：
    #   1) 页面表格内所有单元格文字字体为微软雅黑
    #   2) 字号 10磅
    #   3) 表格保留三列表头"来源""改写后的现场反馈""关键需求"
    #   4) 表格三行内容（即除表头外至少3行数据）
    #   5) 表格边框和底部"核心需求"说明没有明显错位（不超出页面边界、且与表格不重叠）
    s7 = slides[6]
    s7_ok = False
    s7_tables = get_table_shapes(s7)
    slide_w_emu_s7 = prs.slide_width
    slide_h_emu_s7 = prs.slide_height
    expected_headers = ["来源", "改写后的现场反馈", "关键需求"]
    for tshape in s7_tables:
        table = tshape.table
        rows = list(table.rows)
        if len(rows) < 4:  # 表头1行 + 内容3行
            continue
        # (3) 表头三列
        header_cells = [cell.text.strip() for cell in rows[0].cells]
        if len(header_cells) < 3:
            continue
        if not all(h in " ".join(header_cells) for h in expected_headers):
            continue
        # (1)(2) 所有非空 run 字体微软雅黑、字号10磅
        font_ok = True
        for row in rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        fname = get_font_name(run)
                        if fname is None or "微软雅黑" not in fname:
                            font_ok = False
                        sz = get_font_size_pt(run)
                        if sz is None or abs(sz - 10) > 0.5:
                            font_ok = False
        if not font_ok:
            continue
        # (5) 表格不超出页面边界
        if tshape.left is None or tshape.top is None or tshape.width is None or tshape.height is None:
            continue
        if (tshape.left < 0 or tshape.top < 0
                or tshape.left + tshape.width > slide_w_emu_s7
                or tshape.top + tshape.height > slide_h_emu_s7):
            continue
        # 底部"核心需求"说明 shape 不超出页面、不与表格重叠
        core_shape = None
        for shape in s7.shapes:
            if shape is tshape:
                continue
            if shape.has_text_frame and "核心需求" in shape.text_frame.text:
                core_shape = shape
                break
        if core_shape is None:
            continue
        if core_shape.left is None or core_shape.top is None or core_shape.width is None or core_shape.height is None:
            continue
        if (core_shape.left < 0 or core_shape.top < 0
                or core_shape.left + core_shape.width > slide_w_emu_s7
                or core_shape.top + core_shape.height > slide_h_emu_s7):
            continue
        try:
            if check_shape_overlap(tshape, core_shape):
                continue
        except Exception:
            continue
        s7_ok = True
        break
    scores.append((1, "第7页需求表格文字", s7_ok))

    # ---- +5: 第8页三张图尺寸统一 ----
    # 细则要求：
    #   1) 页面中出现三张图片
    #   2) 三张图片均为宽 8cm、高 5cm,允许 ±0.2cm 误差
    #   3) 三张图位置不遮挡"原理 A｜分割""原理 B｜预先作用""组合方向"等正文文字
    s8 = slides[7]
    s8_imgs = get_images_in_slide(s8)
    s8_img_ok = len(s8_imgs) >= 3
    if s8_img_ok:
        for img in s8_imgs[:3]:
            w_cm = emu_to_cm(img.width)
            h_cm = emu_to_cm(img.height)
            if w_cm is None or not (7.8 <= w_cm <= 8.2):
                s8_img_ok = False
                break
            if h_cm is None or not (4.8 <= h_cm <= 5.2):
                s8_img_ok = False
                break
    if s8_img_ok:
        if check_image_covers_text(s8, ["原理 A", "原理A", "分割", "原理 B", "原理B", "预先作用", "组合方向"]):
            s8_img_ok = False
    scores.append((5, "第8页三张图尺寸统一且不遮挡文字", s8_img_ok))

    # ---- +5: 第9页目标说明文字扩充 ----
    # 细则要求：
    #   1) 页面下方说明文本
    #   2) 由原文"目标设定采用区间化与功能化表达,避免绑定真实生产线专属数据。"扩充为约50字的完整句子
    #   3) 仍表达目标设定口径、区间化、功能化、公开交流四层含义
    #   4) 字体微软雅黑、字号20磅
    #   5) 颜色 00A6A6 (浅绿色)
    #   6) 未超出页面底部可视区域
    s9 = slides[8]
    s9_ok = False
    slide_h_emu_s9 = prs.slide_height
    for shape in s9.shapes:
        if not shape.has_text_frame:
            continue
        # (1) 页面下方：shape 顶部位于页面 50% 高度以下
        if shape.top is None or shape.top < slide_h_emu_s9 * 0.5:
            continue
        for para in shape.text_frame.paragraphs:
            ptext = para.text
            # (3) 四层含义：目标设定口径、区间化、功能化、公开交流
            if "目标设定" not in ptext:
                continue
            if "区间" not in ptext:
                continue
            if "功能" not in ptext:
                continue
            if not any(k in ptext for k in ["公开", "交流", "公开交流"]):
                continue
            # (2) 约50字的完整句子(汉字数 ≥ 45,允许"约"的弹性)
            han_count = len(re.findall(r'[一-鿿]', ptext))
            if han_count < 45:
                continue
            # 完整句子：以中文句号/问号/叹号结尾
            stripped = ptext.rstrip()
            if not stripped or stripped[-1] not in '。！？.!?':
                continue
            # (4)(5) 字体微软雅黑、字号20磅、颜色 00A6A6
            fmt_ok = True
            has_run = False
            for run in para.runs:
                if not run.text.strip():
                    continue
                has_run = True
                fname = get_font_name(run)
                if fname is None or "微软雅黑" not in fname:
                    fmt_ok = False
                sz = get_font_size_pt(run)
                if sz is None or abs(sz - 20) > 0.5:
                    fmt_ok = False
                if not color_matches(run, '00A6A6'):
                    fmt_ok = False
            if not (has_run and fmt_ok):
                continue
            # (6) 未超出页面底部可视区域
            if shape.height is None or shape.top + shape.height > slide_h_emu_s9:
                continue
            s9_ok = True
            break
        if s9_ok:
            break
    scores.append((5, "第9页目标说明文字扩充约50字", s9_ok))

    # ---- +3: 第10页图片动态效果 ----
    # 细则要求：
    #   1) 页面中的主要图片对象或图示对象设置"出现"动态效果(进入类动画,presetClass="entr")
    #   2) 动画作用对象位于页面主体图片区
    #   3) 动画不作用于整页所有文字
    #   4) 不导致表格文字消失或错位
    s10 = slides[9]
    slide_w_emu_s10 = prs.slide_width
    slide_h_emu_s10 = prs.slide_height

    slide_xml_s10 = s10._element.xml

    # 收集第10页图片/图示对象与文字/表格对象的 sp id
    img_or_diagram_ids = set()
    text_shape_ids = set()
    table_shape_ids = set()
    main_image_shapes = []
    for shape in s10.shapes:
        try:
            shp_id = shape.shape_id
        except Exception:
            shp_id = id(shape)
        if shape.shape_type == 13:  # PICTURE
            img_or_diagram_ids.add(str(shp_id))
            main_image_shapes.append(shape)
        # 图示对象：组合或 SmartArt/Diagram
        try:
            if shape.shape_type in (6, 12):  # GROUP=6, DIAGRAM/SmartArt 通常作为图片或组合
                img_or_diagram_ids.add(str(shp_id))
        except Exception:
            pass
        if shape.has_text_frame and shape.text_frame.text.strip():
            text_shape_ids.add(str(shp_id))
        if shape.has_table:
            table_shape_ids.add(str(shp_id))

    # (1) 提取"出现"(进入类)动画块及其目标 spTgt id
    #     PPT 中"出现"属于进入类动画,OOXML 通过 presetClass="entr" 标识;
    #     经典"出现"效果 presetID="1" presetClass="entr"。
    entr_targets = set()
    for m in re.finditer(r'<p:par\b[^>]*>.*?</p:par>', slide_xml_s10, re.DOTALL):
        block = m.group(0)
        # 该动画块内含进入类预设
        if 'presetClass="entr"' not in block:
            continue
        for tgt in re.findall(r'<p:spTgt\s+spid="(\d+)"', block):
            entr_targets.add(tgt)
    # 兜底:全文扫描 presetClass="entr" 附近的 spTgt
    if not entr_targets:
        for m in re.finditer(r'presetClass="entr".{0,2000}?<p:spTgt\s+spid="(\d+)"',
                             slide_xml_s10, re.DOTALL):
            entr_targets.add(m.group(1))

    s10_anim_ok = bool(entr_targets)

    # (2) 动画作用对象位于页面主体图片区:
    #     "出现"动画的目标图片中,至少存在一张其中心位于页面主体区域
    if s10_anim_ok:
        in_main_area = False
        for img in main_image_shapes:
            try:
                img_id = str(img.shape_id)
            except Exception:
                img_id = None
            if img_id is None or img_id not in entr_targets:
                continue
            if img.left is None or img.top is None or img.width is None or img.height is None:
                continue
            cx = img.left + img.width / 2
            cy = img.top + img.height / 2
            if (slide_w_emu_s10 * 0.05 <= cx <= slide_w_emu_s10 * 0.95
                    and slide_h_emu_s10 * 0.15 <= cy <= slide_h_emu_s10 * 0.95):
                in_main_area = True
                break
        if not in_main_area:
            s10_anim_ok = False

    # (1)(3)(4) 目标集合校验:
    #     - 必须至少有一个目标是主体图片/图示对象
    #     - 目标集合不能覆盖全部文字 shape
    #     - 不能作用到表格 shape
    if s10_anim_ok:
        # (1) 动画必须作用在主体图片/图示对象上
        if img_or_diagram_ids and not (entr_targets & img_or_diagram_ids):
            s10_anim_ok = False
        # (3) 动画不作用于整页所有文字:目标集合不能覆盖所有文字 shape
        if text_shape_ids and text_shape_ids.issubset(entr_targets):
            s10_anim_ok = False
        # (4) 不作用到表格 shape
        if entr_targets & table_shape_ids:
            s10_anim_ok = False
    scores.append((3, "第10页图片出现动画", s10_anim_ok))

    # ---- +5: 第11页验证页修改 ----
    # ---- +5: 第11页验证页修改 ----
    # 细则要求：
    #   1) "验证结论"右方文字
    #   2) 扩充为50字以上
    #   3) 内容围绕"单层筛面分路""孔型""角度""密封边界优化"展开
    #   4) 字体微软雅黑、字号13
    #   5) 颜色黑色 (000000)
    s11 = slides[10]
    s11_text_ok = False
    conclusion_shape = None
    for shape in s11.shapes:
        if shape.has_text_frame and "验证结论" in shape.text_frame.text:
            conclusion_shape = shape
            break
    if conclusion_shape is not None:
        for shape in s11.shapes:
            if not shape.has_text_frame or shape is conclusion_shape:
                continue
            if shape.left is None:
                continue
            # (1) "右方": shape 左边位于"验证结论"shape 左边之右
            if conclusion_shape.left is None or shape.left < conclusion_shape.left:
                continue
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                # (2) 50字以上(汉字数 ≥ 50)
                if len(re.findall(r'[一-鿿]', t)) < 50:
                    continue
                # (3) 内容关键词:单层筛面分路 / 孔型 / 角度 / 密封边界优化
                if "单层筛面" not in t or "分路" not in t:
                    continue
                if "孔型" not in t:
                    continue
                if "角度" not in t:
                    continue
                if "密封边界" not in t or "优化" not in t:
                    continue
                # (4)(5) 字体微软雅黑、字号13、颜色黑色
                fmt_ok = True
                has_run = False
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    has_run = True
                    fname = get_font_name(run)
                    if fname is None or "微软雅黑" not in fname:
                        fmt_ok = False
                    sz = get_font_size_pt(run)
                    if sz is None or abs(sz - 13) > 0.5:
                        fmt_ok = False
                    # 颜色黑色：放宽为接受固定000000及深色主题色(tx1/dk1等)
                    if not is_run_color_black(run):
                        fmt_ok = False
                if has_run and fmt_ok:
                    s11_text_ok = True
                    break
            if s11_text_ok:
                break
    scores.append((5, "第11页验证结论扩充50字+字号13", s11_text_ok))

    # ---- +3: 第11页轮子动画 ----
    # 细则要求：第11页页面"主要图片对象"设置"轮子"动态效果。
    # PPT"轮子"(Wheel)属于进入类动画,OOXML 以 prstClass="wheel" 或
    # presetID="21" presetClass="entr" 表示。
    s11_xml = s11._element.xml
    wheel_targets = set()
    for m in re.finditer(r'<p:par\b[^>]*>.*?</p:par>', s11_xml, re.DOTALL):
        block = m.group(0)
        block_lower = block.lower()
        is_wheel = (
            'prstclass="wheel"' in block_lower
            or 'filter="wheel"' in block_lower
            or ('presetid="21"' in block_lower and 'presetclass="entr"' in block_lower)
        )
        if not is_wheel:
            continue
        for tgt in re.findall(r'<p:spTgt\s+spid="(\d+)"', block):
            wheel_targets.add(tgt)
    # 兜底:全文扫描 wheel 关键字附近的 spTgt
    if not wheel_targets:
        for m in re.finditer(r'wheel.{0,2000}?<p:spTgt\s+spid="(\d+)"',
                             s11_xml, re.DOTALL | re.IGNORECASE):
            wheel_targets.add(m.group(1))

    s11_anim_ok = bool(wheel_targets)

    if s11_anim_ok:
        # 收集第11页所有图片 shape 的 id,要求轮子动画作用于至少一张图片
        image_ids = set()
        for shape in s11.shapes:
            if shape.shape_type == 13:  # PICTURE
                try:
                    image_ids.add(str(shape.shape_id))
                except Exception:
                    pass

        if not (wheel_targets & image_ids):
            s11_anim_ok = False
    scores.append((3, "第11页图片轮子动画", s11_anim_ok))

    # ---- +1: 第12页试验说明文字 ----
    # 细则要求：
    #   1) 页面底部或下方说明句
    #   2) 完整句子"本试验仅用于工艺趋势判断,不作为真实生产环境的燃烧测试记录。"
    #   3) 字体微软雅黑、字号20磅
    #   4) 颜色深灰色 64748B
    #   5) 文字不超出文本框 (即 text frame 未自动溢出)
    s12 = slides[11]
    s12_target_full = "本试验仅用于工艺趋势判断，不作为真实生产环境的燃烧测试记录。"
    s12_target_alt = "本试验仅用于工艺趋势判断,不作为真实生产环境的燃烧测试记录。"  # 半角逗号兜底
    slide_h_emu_s12 = prs.slide_height
    slide_w_emu_s12 = prs.slide_width
    s12_ok = False

    def _norm(t):
        return t.replace(" ", "").replace("\n", "").replace("\r", "")

    target_norm_a = _norm(s12_target_full)
    target_norm_b = _norm(s12_target_alt)

    for shape in s12.shapes:
        if not shape.has_text_frame:
            continue
        # (1) 页面底部或下方:shape 顶部位于页面 50% 以下
        if shape.top is None or shape.top < slide_h_emu_s12 * 0.5:
            continue
        shape_text_norm = _norm(shape.text_frame.text)
        # (2) 文字完整出现
        if target_norm_a not in shape_text_norm and target_norm_b not in shape_text_norm:
            continue
        # (3)(4) 找到包含"本试验"的 run,校验字体/字号/颜色
        target_run = None
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "本试验" in run.text:
                    target_run = run
                    break
            if target_run:
                break
        if target_run is None:
            continue
        fname = get_font_name(target_run)
        if fname is None or "微软雅黑" not in fname:
            continue
        sz = get_font_size_pt(target_run)
        if sz is None or abs(sz - 20) > 0.5:
            continue
        if not color_matches(target_run, '64748B'):
            continue
        # (5) 文字不超出文本框:shape 完整位于页面内
        if shape.left is None or shape.width is None or shape.height is None:
            continue
        if (shape.left < 0 or shape.left + shape.width > slide_w_emu_s12
                or shape.top + shape.height > slide_h_emu_s12):
            continue
        s12_ok = True
        break
    scores.append((1, "第12页试验说明文字", s12_ok))

    s13 = slides[12]

    # ---- +1: 第13页"设计原则"下方出现五行文字 ----
    # 细则要求:"设计原则"标题下方(不包含"设计原则"所在这一行本身)恰好出现 5 行非空文字。
    # 说明:PPT 中"设计原则\n1. …\n2. …\n3. …\n4. …\n5. …"这种写法在 python-pptx 里
    #      是同一段落内以软换行(<a:br/>)分隔,不会拆分成多个 paragraph;因此按 \n 展开成"行"来数,
    #      并把"设计原则"所在的那一行(可能残留冒号/顿号/空格)整行剔除,只统计其之后的行。
    #      若同框内不足 5 行,再累加位于其下方的其它文本框的非空行。
    def _count_non_empty_lines(text):
        return sum(1 for ln in text.split('\n') if ln.strip())

    def _lines_after_title(full_text, title):
        """返回 full_text 中 title 之后的所有行(丢弃 title 所在整行),按 \n 拆分,过滤空行。"""
        idx = full_text.find(title)
        if idx < 0:
            return []
        tail = full_text[idx + len(title):]
        # 丢弃 title 所在这一行剩余部分:tail 首个 \n 之前的内容都属于标题行
        nl = tail.find('\n')
        if nl < 0:
            return []
        tail = tail[nl + 1:]
        return [ln for ln in tail.split('\n') if ln.strip()]

    s13_five_lines_ok = False
    title_shape_13 = None
    title_tail_lines = 0
    for shape in s13.shapes:
        if not shape.has_text_frame:
            continue
        full_text = shape.text_frame.text
        if "设计原则" not in full_text:
            continue
        title_shape_13 = shape
        title_tail_lines = len(_lines_after_title(full_text, "设计原则"))
        break

    if title_shape_13 is not None:
        if title_tail_lines == 5:
            s13_five_lines_ok = True
        elif title_tail_lines < 5:
            # 同一文本框不足 5 行,再累加位于标题文本框下方的其它文本框的非空行
            below_lines = title_tail_lines
            title_bottom = None
            if title_shape_13.top is not None and title_shape_13.height is not None:
                title_bottom = title_shape_13.top + title_shape_13.height
            for shape in s13.shapes:
                if shape is title_shape_13:
                    continue
                if not shape.has_text_frame:
                    continue
                if shape.top is None:
                    continue
                if title_bottom is not None and shape.top < title_bottom - 100000:
                    continue
                below_lines += _count_non_empty_lines(shape.text_frame.text)
                if below_lines > 5:
                    break
            if below_lines == 5:
                s13_five_lines_ok = True
        # title_tail_lines > 5 视为不满足(超过 5 行)
    scores.append((1, "第13页设计原则下方五行文字", s13_five_lines_ok))

    # ---- +1: 第14页研发指标表格 ----
    # 细则要求：
    #   1) 表格内文字字体微软雅黑、字号20磅、颜色 334155 (深蓝色)
    #   2) 表格包含三列表头："类别""要求""改写后的评价方式"
    #   3) 表格包含五行内容:功能、安全、监测、兼容、成本
    #   4) 底部说明"指标已经转为公开交流口径,不包含原采购金额、供应渠道与真实硬件型号。"
    #   5) 该底部说明:微软雅黑、20磅、颜色 64748B (深灰色)
    s14 = slides[13]
    s14_ok = False
    s14_tables = get_table_shapes(s14)
    expected_headers_14 = ["类别", "要求", "改写后的评价方式"]
    expected_row_keywords_14 = ["功能", "安全", "监测", "兼容", "成本"]
    table_match_ok = False
    for tshape in s14_tables:
        table = tshape.table
        rows = list(table.rows)
        if len(rows) < 6:  # 表头 + 5 行
            continue
        # (2) 三列表头
        header_texts = [cell.text.strip() for cell in rows[0].cells]
        if len(header_texts) < 3:
            continue
        if not all(h in " ".join(header_texts) for h in expected_headers_14):
            continue
        # (3) 五行内容关键词,每个关键词至少出现在一行的第一列(或任一单元格)
        body_rows_text = []
        for row in rows[1:]:
            body_rows_text.append(" ".join(cell.text for cell in row.cells))
        if not all(any(kw in t for t in body_rows_text) for kw in expected_row_keywords_14):
            continue
        # (1) 字体微软雅黑、字号20磅、颜色 334155
        font_ok = True
        for row in rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        fname = get_font_name(run)
                        if fname is None or "微软雅黑" not in fname:
                            font_ok = False
                        sz = get_font_size_pt(run)
                        if sz is None or abs(sz - 20) > 0.5:
                            font_ok = False
                        if not color_matches(run, '334155'):
                            font_ok = False
        if not font_ok:
            continue
        table_match_ok = True
        break

    # (4)(5) 底部说明文字
    footer_ok = False
    footer_norm = "指标已经转为公开交流口径，不包含原采购金额、供应渠道与真实硬件型号。".replace(" ", "")
    footer_norm_alt = footer_norm.replace("，", ",")
    for shape in s14.shapes:
        if not shape.has_text_frame:
            continue
        # 排除表格内的文本帧
        if shape.has_table:
            continue
        shape_text_norm = shape.text_frame.text.replace(" ", "").replace("\n", "").replace("\r", "")
        if footer_norm not in shape_text_norm and footer_norm_alt not in shape_text_norm:
            continue
        # 找含"指标已经"的 run
        target_run = None
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "指标已经" in run.text or "公开交流口径" in run.text:
                    target_run = run
                    break
            if target_run:
                break
        if target_run is None:
            continue
        fname = get_font_name(target_run)
        if fname is None or "微软雅黑" not in fname:
            continue
        sz = get_font_size_pt(target_run)
        if sz is None or abs(sz - 20) > 0.5:
            continue
        if not color_matches(target_run, '64748B'):
            continue
        footer_ok = True
        break

    s14_ok = table_match_ok and footer_ok
    scores.append((1, "第14页研发指标表格", s14_ok))

    # ---- +1: 第15页实施页文字替换 ----
    s15 = slides[14]
    # 细则要求：
    #   1) 表格最后一行"责任角色"列中出现"项目协调负责人：沈亦然"
    #   2) 该单元格字体、字号与同列表格其它行文字一致
    #   3) 表格其他内容仍保持可编辑(即表格仍为真实表格 shape,可遍历单元格)
    s15 = slides[14]
    s15_text = get_slide_text_joined(s15)
    s15_ok = False
    for shape in s15.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        rows = list(table.rows)
        if len(rows) < 2:
            continue
        header = [cell.text.strip() for cell in rows[0].cells]
        if "责任角色" not in header:
            continue
        role_col = header.index("责任角色")
        # (1) 最后一行"责任角色"列文本为"项目协调负责人：沈亦然"
        last_cell = rows[-1].cells[role_col]
        last_text_norm = last_cell.text.strip().replace(":", "：").replace(" ", "")
        if last_text_norm != "项目协调负责人：沈亦然":
            continue
        # (2) 字体、字号与同列其它行一致
        reference_runs = []
        for row_idx in range(1, len(rows) - 1):
            for para in rows[row_idx].cells[role_col].text_frame.paragraphs:
                reference_runs.extend([run for run in para.runs if run.text.strip()])
        target_runs = []
        for para in last_cell.text_frame.paragraphs:
            target_runs.extend([run for run in para.runs if run.text.strip()])
        if not reference_runs or not target_runs:
            continue
        ref_run = reference_runs[0]
        target_run = target_runs[0]
        ref_font = get_font_name(ref_run)
        target_font = get_font_name(target_run)
        ref_size = get_font_size_pt(ref_run)
        target_size = get_font_size_pt(target_run)
        font_consistent = (ref_font == target_font)
        size_consistent = (ref_size is not None and target_size is not None
                           and abs(ref_size - target_size) <= 0.1)
        if not (font_consistent and size_consistent):
            continue
        # (3) 表格其它内容仍可编辑:shape.has_table 已为 True,即真实表格;
        #     再确认表头之外仍有非空单元格(说明内容未被清空/图片化)
        other_has_content = False
        for row_idx in range(1, len(rows) - 1):
            for cell in rows[row_idx].cells:
                if cell.text.strip():
                    other_has_content = True
                    break
            if other_has_content:
                break
        if not other_has_content:
            continue
        s15_ok = True
        break
    scores.append((1, "第15页'项目协调负责人：沈亦然'", s15_ok))

    # ---- +1: 第16页实施一文字字号 ----
    # ---- +1: 第16页实施一文字字号 ----
    # 细则要求：
    #   1) 三句话全部出现:
    #      "长孔筛面兼顾分离效率与堵塞风险"
    #      "让物料在自重与振动下自然前移"
    #      "减少粉料外逸和箱体边缘堆积"
    #   2) 三句字体均为微软雅黑、字号15磅、颜色 334155 (深灰色)
    #   3) 三句话分别位于"1""2""3"三个步骤说明区域内
    s16 = slides[15]
    s16_texts_check = [
        "长孔筛面兼顾分离效率与堵塞风险",
        "让物料在自重与振动下自然前移",
        "减少粉料外逸和箱体边缘堆积",
    ]
    expected_step_labels = [["1", "①", "一"], ["2", "②", "二"], ["3", "③", "三"]]
    s16_ok = True

    sentence_shapes = []
    for target in s16_texts_check:
        # (1) 句子完整出现
        target_norm = target.replace(" ", "")
        shape = None
        target_run = None
        for sh in s16.shapes:
            if not sh.has_text_frame:
                continue
            tf_norm = sh.text_frame.text.replace(" ", "").replace("\n", "")
            if target_norm not in tf_norm:
                continue
            shape = sh
            # 找到包含句首关键词的 run(取前 4 字定位)
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if target[:2] in run.text or target[:3] in run.text or target[:4] in run.text:
                        target_run = run
                        break
                if target_run:
                    break
            break
        if shape is None or target_run is None:
            s16_ok = False
            break
        # (2) 字体微软雅黑、字号15磅、颜色 334155
        fname = get_font_name(target_run)
        if fname is None or "微软雅黑" not in fname:
            s16_ok = False
            break
        sz = get_font_size_pt(target_run)
        if sz is None or abs(sz - 15) > 0.5:
            s16_ok = False
            break
        if not color_matches(target_run, '334155'):
            s16_ok = False
            break
        sentence_shapes.append(shape)

    # (3) 三句话分别位于"1""2""3"三个步骤说明区域内
    if s16_ok:
        if len(sentence_shapes) != 3:
            s16_ok = False
        else:
            # 三个 shape 必须不同
            if len(set(id(sh) for sh in sentence_shapes)) != 3:
                s16_ok = False
        if s16_ok:
            # 每句对应步骤号:同一 shape 内含步骤标签,或附近同一区域含步骤标签
            for idx, sh in enumerate(sentence_shapes):
                labels = expected_step_labels[idx]
                # 先看自身 shape 文本是否包含步骤号
                self_text = sh.text_frame.text.replace(" ", "")
                if any(lbl in self_text for lbl in labels):
                    continue
                # 否则在附近找(矩形重叠 或 同一水平行)
                found_label = False
                for other in s16.shapes:
                    if other is sh or not other.has_text_frame:
                        continue
                    other_text = other.text_frame.text.replace(" ", "").replace("\n", "")
                    if not any(lbl == other_text or lbl in other_text[:2] for lbl in labels):
                        continue
                    try:
                        # 矩形重叠
                        if check_shape_overlap(sh, other):
                            found_label = True
                            break
                        # 同一水平行:垂直中心差距 < max(两 shape 高度)
                        if (sh.top is not None and sh.height is not None
                                and other.top is not None and other.height is not None):
                            sh_cy = sh.top + sh.height / 2
                            ot_cy = other.top + other.height / 2
                            row_tol = max(sh.height, other.height)
                            if abs(sh_cy - ot_cy) <= row_tol:
                                found_label = True
                                break
                        # 同一垂直列:水平中心差距小于 shape 宽度
                        if (sh.left is not None and sh.width is not None
                                and other.left is not None and other.width is not None):
                            sh_cx = sh.left + sh.width / 2
                            ot_cx = other.left + other.width / 2
                            if abs(sh_cx - ot_cx) <= sh.width:
                                found_label = True
                                break
                    except Exception:
                        pass
                if not found_label:
                    s16_ok = False
                    break
    scores.append((1, "第16页三句话微软雅黑15磅", s16_ok))

    # ---- +1: 第17页实施二表格与动画 ----
    # 细则要求：
    #   1) 页面表格内文字字体微软雅黑、字号10磅
    #   2) 表格仍包含"检查项""测试结果""开合便利性""内部可视度""维护清洁"等文字
    #   3) 文字颜色深灰色 334155
    s17 = slides[16]
    s17_table_ok = False
    s17_required_keywords = ["检查项", "测试结果", "开合便利性", "内部可视度", "维护清洁"]
    for tshape in get_table_shapes(s17):
        table = tshape.table
        # (2) 表格内必须包含所有指定关键词
        all_cell_text = []
        for row in table.rows:
            for cell in row.cells:
                all_cell_text.append(cell.text)
        joined = " ".join(all_cell_text)
        if not all(kw in joined for kw in s17_required_keywords):
            continue
        # (1)(3) 字体微软雅黑、字号10磅、颜色 334155
        font_ok = True
        for row in table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        fname = get_font_name(run)
                        if fname is None or "微软雅黑" not in fname:
                            font_ok = False
                        sz = get_font_size_pt(run)
                        if sz is None or abs(sz - 10) > 0.5:
                            font_ok = False
                        if not color_matches(run, '334155'):
                            font_ok = False
        if not font_ok:
            continue
        s17_table_ok = True
        break
    scores.append((1, "第17页表格文字", s17_table_ok))

    # ---- +3: 第17页渐变动画 ----
    # 细则要求：第17页页面"主要图片对象"设置"渐变"动态效果。
    # PPT"渐变"(Fade)属于进入类动画,OOXML 以 filter/prstClass="fade"
    # 或 presetID="10" presetClass="entr" 标识。
    s17_xml = s17._element.xml
    fade_targets = set()
    for m in re.finditer(r'<p:par\b.*?</p:par>', s17_xml, re.DOTALL):
        block = m.group(0)
        bl = block.lower()
        is_fade = (
            'prstclass="fade"' in bl
            or 'filter="fade"' in bl
            or ('presetid="10"' in bl and 'presetclass="entr"' in bl)
        )
        if not is_fade:
            continue
        for tgt in re.findall(r'<p:spTgt\s+spid="(\d+)"', block):
            fade_targets.add(tgt)
    # 兜底:全文扫描 fade 关键字附近的 spTgt
    if not fade_targets:
        for m in re.finditer(r'fade.{0,2000}?<p:spTgt\s+spid="(\d+)"',
                             s17_xml, re.DOTALL | re.IGNORECASE):
            fade_targets.add(m.group(1))

    s17_anim_ok = bool(fade_targets)

    if s17_anim_ok:
        # 收集第17页所有图片 shape 的 id,要求渐变动画作用于至少一张图片
        s17_image_ids = set()
        for shape in s17.shapes:
            if shape.shape_type == 13:  # PICTURE
                try:
                    s17_image_ids.add(str(shape.shape_id))
                except Exception:
                    pass

        if not (fade_targets & s17_image_ids):
            s17_anim_ok = False
    scores.append((3, "第17页图片渐变动画", s17_anim_ok))

    # ---- +3: 第18页动画 ----
    # 细则要求：
    #   1) 页面主要图片对象设置溶解(dissolve)动态效果
    #   2) 图片位置仍位于页面主体区域
    #   3) 不遮挡"画面覆盖""信号独立""巡检记录"三组文字
    s18 = slides[17]
    slide_w_emu_s18 = prs.slide_width
    slide_h_emu_s18 = prs.slide_height

    s18_anim_ok = check_animation_effect(s18, '溶解')

    # (1) 解析 slide XML 中溶解(dissolve)动画的 spTgt 目标 id,且至少作用于一张图片
    s18_image_shapes = [sh for sh in s18.shapes if sh.shape_type == 13]
    if s18_anim_ok:
        s18_xml = s18._element.xml
        dissolve_targets = set()
        for m in re.finditer(r'<p:par\b[^>]*>.*?</p:par>', s18_xml, re.DOTALL):
            block = m.group(0)
            if 'dissolve' not in block.lower():
                continue
            for tgt in re.findall(r'<p:spTgt\s+spid="(\d+)"', block):
                dissolve_targets.add(tgt)
        if not dissolve_targets:
            for m in re.finditer(r'dissolve.{0,2000}?<p:spTgt\s+spid="(\d+)"', s18_xml, re.DOTALL | re.IGNORECASE):
                dissolve_targets.add(m.group(1))
        image_ids = set()
        for img in s18_image_shapes:
            try:
                image_ids.add(str(img.shape_id))
            except Exception:
                pass
        if not dissolve_targets or not (dissolve_targets & image_ids):
            s18_anim_ok = False

    # (2) 图片仍位于页面主体区域:至少一张图片中心位于 (水平 5%~95%, 垂直 15%~95%)
    if s18_anim_ok:
        in_main_area = False
        for img in s18_image_shapes:
            if img.left is None or img.top is None or img.width is None or img.height is None:
                continue
            cx = img.left + img.width / 2
            cy = img.top + img.height / 2
            if (slide_w_emu_s18 * 0.05 <= cx <= slide_w_emu_s18 * 0.95
                    and slide_h_emu_s18 * 0.15 <= cy <= slide_h_emu_s18 * 0.95):
                in_main_area = True
                break
        if not in_main_area:
            s18_anim_ok = False

    # (3) 三组文字全部出现,且图片不遮挡这些文字
    s18_check_texts = ["画面覆盖", "信号独立", "巡检记录"]
    s18_text_check = all(slide_contains_text(s18, kw) for kw in s18_check_texts)
    s18_no_cover = not check_image_covers_text(s18, s18_check_texts)
    scores.append((3, "第18页溶解动画+不遮挡文字", s18_anim_ok and s18_text_check and s18_no_cover))

    # ---- +1: 第19页Excel数据表格 ----
    # 细则要求：
    #   1) 页面上方表格(内容来自上传Excel的 Sheet1!A1:D6,即 6行 × 4列)
    #   2) 表头4列:"评估项目""方案A:紧凑型产线""方案B:标准化产线""方案C:智能柔性产线"
    #   3) 行项目5项:"单线产能""工艺配置""维护空间""运行稳定性""综合关注点"
    #   4) 表格整体高度 10.59cm,宽度 24.07cm,允许 ±0.2cm 误差
    #   5) 表格文字可编辑(真实表格 shape)
    s19 = slides[18]
    s19_table_ok = False
    s19_tables = get_table_shapes(s19)
    slide_h_emu_s19 = prs.slide_height
    expected_headers_19 = [
        "评估项目",
        "方案A：紧凑型产线",
        "方案B：标准化产线",
        "方案C：智能柔性产线",
    ]
    expected_rows_19 = ["单线产能", "工艺配置", "维护空间", "运行稳定性", "综合关注点"]
    for tshape in s19_tables:
        # (5) 真实表格:可编辑
        table = tshape.table
        rows = list(table.rows)
        # (1) Sheet1!A1:D6 → 6 行 × 4 列
        if len(rows) < 6:
            continue
        # 取首行 4 列校验表头
        header_cells = [cell.text.strip().replace(":", "：").replace(" ", "") for cell in rows[0].cells]
        if len(header_cells) < 4:
            continue
        # (2) 四个表头
        header_joined = "|".join(header_cells)
        if not all(h.replace(" ", "") in header_joined for h in expected_headers_19):
            continue
        # (3) 行项目 5 项(在第 1 列前 5 行)
        row_labels = []
        for r_idx in range(1, min(6, len(rows))):
            last_text_norm = rows[r_idx].cells[0].text.strip().replace(" ", "")
            row_labels.append(last_text_norm)
        if not all(any(label in rl for rl in row_labels) for label in expected_rows_19):
            continue
        # (4) 尺寸:24.07 × 10.59 cm,±0.2 cm
        w_cm = emu_to_cm(tshape.width)
        h_cm = emu_to_cm(tshape.height)
        if w_cm is None or h_cm is None:
            continue
        if not (abs(w_cm - 24.07) <= 0.2 and abs(h_cm - 10.59) <= 0.2):
            continue
        # (1) 页面上方:表格垂直中心位于页面 50% 之上
        if tshape.top is None or tshape.height is None:
            continue
        center = tshape.top + tshape.height / 2
        if center > slide_h_emu_s19 * 0.5:
            continue
        s19_table_ok = True
        break
    s19_text = get_slide_text_joined(s19)
    scores.append((1, "第19页Excel表格尺寸和表头", s19_table_ok))

    # ---- +1: 第19页Excel表格具体数据 ----
    # 细则要求：
    #   1) 表格中出现下列 Excel 数据文本:
    #      "约7,200件/小时""约9,600件/小时"
    #      "在线检测+分段除尘""常规传送+集中抽吸"
    #      "视觉识别+闭环调节+分级回收"
    #      "粉尘积聚、温升预警与盲区巡检联动"
    #   2) 数据顺序与 Excel 表格从左到右、从上到下顺序一致
    required_cells = [
        "约7,200件/小时",
        "约9,600件/小时",
        "在线检测+分段除尘",
        "常规传送+集中抽吸",
        "视觉识别+闭环调节+分级回收",
        "粉尘积聚、温升预警与盲区巡检联动",
    ]
    def _norm_cell(t):
        return t.replace(" ", "").replace("\n", "").replace("\r", "").replace(",", "，")

    s19_data_ok = False
    yield_pattern = re.compile(r'约?[\d,]+件/小时')
    allowed_yields = {"约7,200件/小时", "约9,600件/小时"}
    allowed_yields_norm = {_norm_cell(x) for x in allowed_yields}
    for tshape in s19_tables:
        # 按从上到下、从左到右展开所有单元格文本
        flat = []
        for row in tshape.table.rows:
            for cell in row.cells:
                flat.append(_norm_cell(cell.text))
        # (1) 全部出现
        norm_required = [_norm_cell(c) for c in required_cells]
        if not all(any(req in t for t in flat) for req in norm_required):
            continue
        # (2) 顺序一致:按 flat 出现顺序与 required 顺序匹配
        order_idx = -1
        order_ok = True
        for req in norm_required:
            found = -1
            for j in range(order_idx + 1, len(flat)):
                if req in flat[j]:
                    found = j
                    break
            if found < 0:
                order_ok = False
                break
            order_idx = found
        if not order_ok:
            continue
        # (3) 表格中不得出现细则未列举的"件/小时"产能数据
        #     (Excel 源数据里方案 C 单线产能为空,若填入"约12,800件/小时"等值即判失败)
        raw_text = "".join(cell.text for row in tshape.table.rows for cell in row.cells)
        yields_found = {_norm_cell(m) for m in yield_pattern.findall(raw_text)}
        if not yields_found.issubset(allowed_yields_norm):
            continue
        s19_data_ok = True
        break
    scores.append((1, "第19页Excel表格具体数据", s19_data_ok))

    # ---- +3: 第20页风险链条圆形颜色 ----
    # 细则要求：
    #   1) "燃料侧"上方圆形 — 填充红色 + 外框线条红色
    #   2) "发现侧"上方圆形 — 填充黄色 + 外框线条黄色
    #   3) 两个圆形直径与页面其他同类圆形一致
    #   4) 圆形不遮挡下方标题文字
    s20 = slides[19]
    s20_ok = False

    def _shape_prst(shape):
        match = re.search(r'<a:prstGeom prst="([^"]+)"', shape._element.xml)
        return match.group(1) if match else None

    def _shape_fill_color(shape):
        sp_pr_match = re.search(r'<p:spPr[^>]*>(.*?)</p:spPr>', shape._element.xml, re.DOTALL)
        if not sp_pr_match:
            return None
        sp_pr = sp_pr_match.group(1)
        before_line = sp_pr.split('<a:ln', 1)[0]
        fill_match = re.search(r'<a:solidFill>\s*<a:srgbClr val="([^"]+)"', before_line)
        return fill_match.group(1).upper() if fill_match else None

    def _shape_line_color(shape):
        line_match = re.search(r'<a:ln[^>]*>.*?<a:solidFill>\s*<a:srgbClr val="([^"]+)"', shape._element.xml, re.DOTALL)
        return line_match.group(1).upper() if line_match else None

    def _find_text_shape(slide, text):
        for shape in slide.shapes:
            if shape.has_text_frame and text in shape.text_frame.text:
                return shape
        return None

    def _find_circle_above(slide, text_shape):
        if text_shape is None:
            return None
        text_center = text_shape.left + text_shape.width / 2
        candidates = []
        for shape in slide.shapes:
            if _shape_prst(shape) != 'ellipse':
                continue
            if abs(shape.width - shape.height) > 1000:
                continue
            circle_center = shape.left + shape.width / 2
            if shape.top < text_shape.top and abs(circle_center - text_center) <= text_shape.width / 2:
                candidates.append(shape)
        if not candidates:
            return None
        candidates.sort(key=lambda shape: abs((shape.left + shape.width / 2) - text_center))
        return candidates[0]

    def _is_red(hex_color):
        if not hex_color:
            return False
        try:
            r = int(hex_color[0:2], 16); g = int(hex_color[2:4], 16); b = int(hex_color[4:6], 16)
        except Exception:
            return False
        # 红色:R 显著大于 G、B,且整体偏红
        return r >= 150 and r > g + 40 and r > b + 40

    def _is_yellow(hex_color):
        if not hex_color:
            return False
        try:
            r = int(hex_color[0:2], 16); g = int(hex_color[2:4], 16); b = int(hex_color[4:6], 16)
        except Exception:
            return False
        # 黄色:R 和 G 都高且接近,B 显著低
        return r >= 180 and g >= 140 and abs(r - g) <= 80 and b < min(r, g) - 60

    # 找到页面上所有"同类圆形"(prstGeom=ellipse 且接近正圆),用于直径一致校验
    all_circles = []
    for shape in s20.shapes:
        if _shape_prst(shape) != 'ellipse':
            continue
        if shape.width is None or shape.height is None:
            continue
        if abs(shape.width - shape.height) > 1000:
            continue
        all_circles.append(shape)

    fuel_text = _find_text_shape(s20, "燃料侧")
    find_text = _find_text_shape(s20, "发现侧")
    fuel_circle = _find_circle_above(s20, fuel_text)
    find_circle = _find_circle_above(s20, find_text)
    if fuel_circle is not None and find_circle is not None:
        # (1) 燃料侧:填充和外框线均红色
        fuel_fill = _shape_fill_color(fuel_circle)
        fuel_line = _shape_line_color(fuel_circle)
        fuel_color_ok = _is_red(fuel_fill) and _is_red(fuel_line)
        # (2) 发现侧:填充和外框线均黄色
        find_fill = _shape_fill_color(find_circle)
        find_line = _shape_line_color(find_circle)
        find_color_ok = _is_yellow(find_fill) and _is_yellow(find_line)
        # (3) 直径与页面其他同类圆形一致
        # 收集除两个目标圆以外的"其他同类圆形"
        other_circles = [c for c in all_circles
                         if c is not fuel_circle and c is not find_circle]
        diameter_consistent = True
        if other_circles:
            ref_w = other_circles[0].width
            ref_h = other_circles[0].height
            for c in (fuel_circle, find_circle):
                if abs(c.width - ref_w) > 50000 or abs(c.height - ref_h) > 50000:  # 容差约 0.5cm
                    diameter_consistent = False
                    break
        # 两个目标圆彼此也要一致
        if abs(fuel_circle.width - find_circle.width) > 50000 or abs(fuel_circle.height - find_circle.height) > 50000:
            diameter_consistent = False
        # (4) 圆形不遮挡下方标题文字
        no_cover = not check_shape_overlap(fuel_circle, fuel_text) and not check_shape_overlap(find_circle, find_text)
        s20_ok = fuel_color_ok and find_color_ok and diameter_consistent and no_cover
    scores.append((3, "第20页燃料侧红色/发现侧黄色圆形", s20_ok))

    # ---- +1: 第21页巩固措施说明 ----
    # 细则要求：
    #   1) 页面下方说明句完整出现:
    #      "下一步：根据不同线体空间差异，形成可调整的安装包和巡检说明卡。"
    #   2) 字体微软雅黑、字号20磅
    #   3) 颜色深蓝色 0A1F33
    #   4) 位于页面底部说明区域
    #   5) 不超出页面边界
    s21 = slides[20]
    slide_w_emu_s21 = prs.slide_width
    slide_h_emu_s21 = prs.slide_height
    target_full_s21 = "下一步：根据不同线体空间差异，形成可调整的安装包和巡检说明卡。"
    target_norm_full_s21 = target_full_s21.replace(" ", "").replace("\n", "")
    target_norm_full_s21_alt = target_norm_full_s21.replace("：", ":").replace("，", ",")
    s21_ok = False
    for shape in s21.shapes:
        if not shape.has_text_frame:
            continue
        # (1) 完整文字出现
        shape_text_norm = shape.text_frame.text.replace(" ", "").replace("\n", "").replace("\r", "")
        if (target_norm_full_s21 not in shape_text_norm
                and target_norm_full_s21_alt not in shape_text_norm.replace("：", ":").replace("，", ",")):
            continue
        # (4) 页面底部说明区域:shape 顶部位于页面 60% 以下
        if shape.top is None or shape.top < slide_h_emu_s21 * 0.6:
            continue
        # (5) 不超出页面边界
        if shape.left is None or shape.width is None or shape.height is None:
            continue
        if (shape.left < 0 or shape.top < 0
                or shape.left + shape.width > slide_w_emu_s21
                or shape.top + shape.height > slide_h_emu_s21):
            continue
        # (2)(3) 字体微软雅黑、字号20磅、颜色 0A1F33
        target_run = None
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "下一步" in run.text:
                    target_run = run
                    break
            if target_run:
                break
        if target_run is None:
            continue
        fname = get_font_name(target_run)
        if fname is None or "微软雅黑" not in fname:
            continue
        sz = get_font_size_pt(target_run)
        if sz is None or abs(sz - 20) > 0.5:
            continue
        if not color_matches(target_run, '0A1F33'):
            continue
        s21_ok = True
        break
    scores.append((1, "第21页巩固措施说明", s21_ok))

    # ---- +1: 第22页总结与展望文字 ----
    # 细则要求：
    #   1) 标题"总结与展望" — 微软雅黑40磅、颜色 FFFFFF
    #   2) 标题位于页面左上或中上标题区域
    #   3) 页面下方完整句:
    #      "后续方向：结合图像识别与清理记录，实现异常趋势提醒与点检数据沉淀。"
    #   4) 后续方向句 — 微软雅黑20磅、颜色 FFFFFF
    #   5) 不遮挡"小改造""强可视""易维护"三组内容
    s22 = slides[21]
    slide_w_emu_s22 = prs.slide_width
    slide_h_emu_s22 = prs.slide_height

    # (1)(2) 标题"总结与展望"
    title_ok_s22 = False
    title_shape_s22 = None
    for shape in s22.shapes:
        if not shape.has_text_frame:
            continue
        if "总结与展望" not in shape.text_frame.text:
            continue
        # 位置:左上或中上 → top 在页面上半部、left 在 0~75%
        if shape.top is None or shape.left is None:
            continue
        if shape.top > slide_h_emu_s22 * 0.5:
            continue
        if shape.left < 0 or shape.left > slide_w_emu_s22 * 0.75:
            continue
        # 字体微软雅黑、字号40磅、颜色 FFFFFF
        title_run = None
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "总结与展望" in run.text:
                    title_run = run
                    break
            if title_run:
                break
        if title_run is None:
            continue
        fname = get_font_name(title_run)
        if fname is None or "微软雅黑" not in fname:
            continue
        sz = get_font_size_pt(title_run)
        if sz is None or abs(sz - 40) > 0.5:
            continue
        if not color_matches(title_run, 'FFFFFF'):
            continue
        title_ok_s22 = True
        title_shape_s22 = shape
        break

    # (3)(4) "后续方向"完整句
    s22_target_full = "后续方向：结合图像识别与清理记录，实现异常趋势提醒与点检数据沉淀。"
    s22_target_norm = s22_target_full.replace(" ", "").replace("\n", "")
    s22_target_norm_alt = s22_target_norm.replace("：", ":").replace("，", ",")
    detail_ok_s22 = False
    detail_shape_s22 = None
    for shape in s22.shapes:
        if not shape.has_text_frame:
            continue
        shape_text_norm = shape.text_frame.text.replace(" ", "").replace("\n", "").replace("\r", "")
        if (s22_target_norm not in shape_text_norm
                and s22_target_norm_alt not in shape_text_norm.replace("：", ":").replace("，", ",")):
            continue
        # 位置:页面下方
        if shape.top is None or shape.top < slide_h_emu_s22 * 0.5:
            continue
        # 微软雅黑、20磅、FFFFFF
        target_run = None
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "后续方向" in run.text:
                    target_run = run
                    break
            if target_run:
                break
        if target_run is None:
            continue
        fname = get_font_name(target_run)
        if fname is None or "微软雅黑" not in fname:
            continue
        sz = get_font_size_pt(target_run)
        if sz is None or abs(sz - 20) > 0.5:
            continue
        if not color_matches(target_run, 'FFFFFF'):
            continue
        detail_ok_s22 = True
        detail_shape_s22 = shape
        break

    # (5) 不遮挡"小改造""强可视""易维护"三组内容
    no_cover_ok_s22 = True
    s22_content_keywords = ["小改造", "强可视", "易维护"]
    content_shapes = []
    for shape in s22.shapes:
        if not shape.has_text_frame:
            continue
        if any(kw in shape.text_frame.text for kw in s22_content_keywords):
            content_shapes.append(shape)
    for title_or_detail in (title_shape_s22, detail_shape_s22):
        if title_or_detail is None:
            continue
        for cs in content_shapes:
            if cs is title_or_detail:
                continue
            try:
                if check_shape_overlap(title_or_detail, cs):
                    no_cover_ok_s22 = False
                    break
            except Exception:
                pass
        if not no_cover_ok_s22:
            break

    s22_total_ok = title_ok_s22 and detail_ok_s22 and no_cover_ok_s22
    scores.append((1, "第22页总结与展望文字", s22_total_ok))

    # ---- +1: 第23页感谢页标题 ----
    # 细则要求：
    #   1) "感谢聆听" — 微软雅黑40磅、颜色 FFFFFF
    #   2) "感谢聆听" 位于页面中部或中上区域
    #   3) "请批评指正"完整出现
    #   4) "持续改进，让隐蔽风险看得见、管得住。"完整出现
    #   5) 文字可编辑(即文本以真实 textframe/run 形式存在,而非图片)
    s23 = slides[22]
    slide_h_emu_s23 = prs.slide_height

    # (1)(2) "感谢聆听"
    s23_title_ok = False
    for shape in s23.shapes:
        if not shape.has_text_frame:
            continue
        if "感谢聆听" not in shape.text_frame.text:
            continue
        # 位置:中部或中上 → 垂直中心位于 20% ~ 70% 高度区间
        if shape.top is None or shape.height is None:
            continue
        cy = shape.top + shape.height / 2
        if not (slide_h_emu_s23 * 0.20 <= cy <= slide_h_emu_s23 * 0.70):
            continue
        # 字体微软雅黑、字号40磅、颜色 FFFFFF
        title_run = None
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "感谢聆听" in run.text:
                    title_run = run
                    break
            if title_run:
                break
        if title_run is None:
            continue
        fname = get_font_name(title_run)
        if fname is None or "微软雅黑" not in fname:
            continue
        sz = get_font_size_pt(title_run)
        if sz is None or abs(sz - 40) > 0.5:
            continue
        if not color_matches(title_run, 'FFFFFF'):
            continue
        s23_title_ok = True
        break

    # (3) "请批评指正" 完整出现 + 可编辑
    s23_extra1_ok = False
    for shape in s23.shapes:
        if not shape.has_text_frame:
            continue
        # 必须以 run 形式出现,确认可编辑
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "请批评指正" in run.text:
                    s23_extra1_ok = True
                    break
            if s23_extra1_ok:
                break
        if s23_extra1_ok:
            break

    # (4) "持续改进，让隐蔽风险看得见、管得住。"完整出现 + 可编辑
    s23_extra2_target = "持续改进，让隐蔽风险看得见、管得住。"
    s23_extra2_norm = s23_extra2_target.replace(" ", "")
    s23_extra2_norm_alt = s23_extra2_norm.replace("，", ",")
    s23_extra2_ok = False
    for shape in s23.shapes:
        if not shape.has_text_frame:
            continue
        shape_text_norm = shape.text_frame.text.replace(" ", "").replace("\n", "").replace("\r", "")
        if (s23_extra2_norm not in shape_text_norm
                and s23_extra2_norm_alt not in shape_text_norm.replace("，", ",")):
            continue
        # 必须以 run 形式承载(可编辑)
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "持续改进" in run.text or "隐蔽风险" in run.text or "管得住" in run.text:
                    s23_extra2_ok = True
                    break
            if s23_extra2_ok:
                break
        if s23_extra2_ok:
            break

    scores.append((1, "第23页感谢页标题", s23_title_ok and s23_extra1_ok and s23_extra2_ok))

    # ============ 扣分项 ============
    deductions = []

    # -3: 除第19页表格外,其余表格内文字对齐方式非左对齐
    # 细则要求:第19页表格不在检查范围内;其余页面只要有任一表格的任一单元格段落显式设为非左对齐,即触发扣分。
    non_left_align = False
    for i, slide in enumerate(slides):
        if i == 18:  # 第19页(1-indexed=19, 0-indexed=18)跳过
            continue
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        if not para.text.strip():
                            continue
                        # None = 继承默认 = 左对齐,不触发扣分
                        # 显式设为非 LEFT 才视为"非左对齐"
                        if para.alignment is not None and para.alignment != PP_ALIGN.LEFT:
                            non_left_align = True
                            break
                    if non_left_align:
                        break
                if non_left_align:
                    break
            if non_left_align:
                break
        if non_left_align:
            break
    deductions.append((-3, "除第19页外表格文字非左对齐", non_left_align))

    # -3: 第8页文字被图片遮挡
    # 细则要求:第8页任一图片矩形与任一非空文本 shape 矩形相交,即视为"文字被图片遮挡"。
    s8_overlap = False
    s8_imgs_for_deduct = get_images_in_slide(s8)
    if s8_imgs_for_deduct:
        for img in s8_imgs_for_deduct:
            for shape in s8.shapes:
                if shape is img:
                    continue
                if not shape.has_text_frame:
                    continue
                if not shape.text_frame.text.strip():
                    continue
                try:
                    if check_shape_overlap(img, shape):
                        s8_overlap = True
                        break
                except Exception:
                    pass
            if s8_overlap:
                break
    deductions.append((-3, "第8页文字被图片遮挡", s8_overlap))

    # -3: 第20页文本框出现多余黄色背景 —— 已删除
    # 原因:与 +3 "第20页燃料侧红色/发现侧黄色圆形" 自相矛盾——
    #   题目本身要求"发现侧"上方圆形必须填充黄色,而该圆 shape 自带 textframe,
    #   会同时被本扣分项判中,造成同一对象既加 3 分又扣 3 分。

    # -1: 第2页没有出现标题文本"目录"或没有出现英文文本"CONTENTS"
    # 细则要求:两者任一缺失即触发 -1 扣分。
    s2_full_text = get_slide_text_joined(s2)
    s2_has_mulu = "目录" in s2_full_text
    s2_has_contents = "CONTENTS" in s2_full_text.upper()
    s2_missing = (not s2_has_mulu) or (not s2_has_contents)
    deductions.append((-1, "第2页缺少'目录'或'CONTENTS'", s2_missing))

    # -1: 第6页仍出现标题文本"目录"或英文文本"CONTENTS"
    # 细则要求:第6页中出现"目录"或"CONTENTS"任一,即触发扣分。
    s6_text = get_slide_text_joined(s6)
    s6_has_mulu = "目录" in s6_text
    s6_has_contents = "CONTENTS" in s6_text.upper()
    s6_deduct = s6_has_mulu or s6_has_contents
    deductions.append((-1, "第6页仍出现'目录'或'CONTENTS'", s6_deduct))

    # -1: 第1页没有出现标题文本"高速成型线碎屑分离与状态监测装置改进汇报"
    # 细则要求:第1页若缺少该标题原文,即触发 -1 扣分。
    s1 = slides[0]
    s1_target = "高速成型线碎屑分离与状态监测装置改进汇报"
    s1_target_norm = s1_target.replace(" ", "")
    s1_missing = True
    for shape in s1.shapes:
        if not shape.has_text_frame:
            continue
        shape_text_norm = shape.text_frame.text.replace(" ", "").replace("\n", "").replace("\r", "")
        if s1_target_norm in shape_text_norm:
            s1_missing = False
            break
    deductions.append((-1, "第1页缺少标题", s1_missing))

    # -3: 第4页三张新增图片中任意一张覆盖"粉状碎屑""残留热斑""封闭堆积"三个标题文本之一
    # 细则要求:任一图片矩形与含上述三标题之一的文本 shape 矩形相交,即触发 -3 扣分。
    s4_cover = check_image_covers_text(s4, ["粉状碎屑", "残留热斑", "封闭堆积"])
    deductions.append((-3, "第4页图片覆盖标题文字", s4_cover))

    # -1: 第11页没有出现文本"05 验证",或第11页出现文本"04 验证"
    # 细则要求:任一条件成立即触发扣分。
    s11_text = get_slide_text_joined(s11)
    s11_text_norm = s11_text.replace(" ", "")
    has_05_verify = "05验证" in s11_text_norm or "05 验证" in s11_text
    has_04_verify = "04验证" in s11_text_norm or "04 验证" in s11_text
    s11_deduct = (not has_05_verify) or has_04_verify
    deductions.append((-1, "第11页缺'05 验证'或出现'04 验证'", s11_deduct))

    # -1: 第15页表格最后一行"责任角色"列出现文本"项目协调：沈亦然"
    # 细则要求:严格定位"最后一行的责任角色列"单元格,且其文本恰为"项目协调：沈亦然"(非"项目协调负责人：沈亦然")。
    s15_wrong = False
    for shape in s15.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        rows = list(table.rows)
        if len(rows) < 2:
            continue
        header = [cell.text.strip() for cell in rows[0].cells]
        if "责任角色" not in header:
            continue
        role_col = header.index("责任角色")
        last_cell = rows[-1].cells[role_col]
        last_text_norm = last_cell.text.strip().replace(":", "：").replace(" ", "")
        if last_text_norm == "项目协调：沈亦然":
            s15_wrong = True
        break
    deductions.append((-1, "第15页出现'项目协调：沈亦然'(非负责人)", s15_wrong))

    # -3: 第16页缺少三句话中任一句
    s16_text_full = get_slide_text_joined(s16)
    # -3: 第16页没有出现三句话中的任意一句即触发
    s16_text_full = get_slide_text_joined(s16).replace(" ", "")
    s16_missing = False
    for t in s16_texts_check:
        if t.replace(" ", "") not in s16_text_full:
            s16_missing = True
            break
    deductions.append((-3, "第16页缺少三句话之一", s16_missing))

    # -1: 第17页表格内任意文字字号 >12磅 或 <8磅
    s17_font_bad = False
    for shape in s17.shapes:
        if not shape.has_table:
            continue
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        sz = get_font_size_pt(run)
                        if sz is not None and (sz > 12 or sz < 8):
                            s17_font_bad = True
                            break
                    if s17_font_bad:
                        break
                if s17_font_bad:
                    break
            if s17_font_bad:
                break
        if s17_font_bad:
            break
    deductions.append((-1, "第17页表格字号超范围", s17_font_bad))

    # -5: 第19页表格没有出现表头四项中的任一项即触发
    # 细则要求(全部都需出现,缺一即扣):
    #   "评估项目""方案A：紧凑型产线""方案B：标准化产线""方案C：智能柔性产线"
    s19_required_headers = [
        "评估项目",
        "方案A：紧凑型产线",
        "方案B：标准化产线",
        "方案C：智能柔性产线",
    ]
    s19_text_norm = s19_text.replace(" ", "").replace(":", "：")
    s19_header_missing = False
    for h in s19_required_headers:
        if h.replace(" ", "") not in s19_text_norm:
            s19_header_missing = True
            break
    deductions.append((-5, "第19页表格缺表头", s19_header_missing))

    # -5: 第19页表格没有出现以下数据文本任一项即触发
    #   "视觉识别+闭环调节+分级回收" / "粉尘积聚、温升预警与盲区巡检联动"
    def _norm19(t):
        return t.replace(" ", "").replace("\n", "").replace(",", "，")
    s19_required_data = [
        "视觉识别+闭环调节+分级回收",
        "粉尘积聚、温升预警与盲区巡检联动",
    ]
    s19_text_norm_for_data = _norm19(s19_text)
    s19_data_deduct = False
    for d in s19_required_data:
        if _norm19(d) not in s19_text_norm_for_data:
            s19_data_deduct = True
            break
    deductions.append((-5, "第19页表格缺关键数据", s19_data_deduct))

    # -3: 第19页表格宽度 <23.8cm 或 >24.3cm,或高度 <10.3cm 或 >10.9cm
    s19_size_bad = False
    if not s19_tables:
        s19_size_bad = True
    else:
        for tshape in s19_tables:
            w_cm = emu_to_cm(tshape.width)
            h_cm = emu_to_cm(tshape.height)
            if w_cm is None or h_cm is None:
                s19_size_bad = True
                break
            if w_cm < 23.8 or w_cm > 24.3:
                s19_size_bad = True
                break
            if h_cm < 10.3 or h_cm > 10.9:
                s19_size_bad = True
                break
    deductions.append((-3, "第19页表格尺寸不符", s19_size_bad))

    # -3: 第21页没有出现完整文本"下一步：根据不同线体空间差异，形成可调整的安装包和巡检说明卡。"
    s21_text_full = get_slide_text_joined(s21).replace(" ", "").replace("\n", "")
    s21_target_full = "下一步：根据不同线体空间差异，形成可调整的安装包和巡检说明卡。".replace(" ", "")
    s21_target_alt = s21_target_full.replace("：", ":").replace("，", ",")
    s21_text_alt = s21_text_full.replace("：", ":").replace("，", ",")
    s21_missing = (s21_target_full not in s21_text_full) and (s21_target_alt not in s21_text_alt)
    deductions.append((-3, "第21页缺少'下一步'文本", s21_missing))

    # -1: 第22页标题"总结与展望"字号 <36磅 或 >44磅
    s22_font_bad = False
    runs22 = find_text_runs_in_slide(s22, "总结与展望")
    if not runs22:
        s22_font_bad = True
    else:
        sz = runs22[0]['font_size_pt']
        if sz is None or sz < 36 or sz > 44:
            s22_font_bad = True
    deductions.append((-1, "第22页标题字号不符(36-44磅)", s22_font_bad))

    # -1: 第23页标题"感谢聆听"字号 <36磅 或 >44磅
    s23_font_bad = False
    runs23 = find_text_runs_in_slide(s23, "感谢聆听")
    if not runs23:
        s23_font_bad = True
    else:
        sz = runs23[0]['font_size_pt']
        if sz is None or sz < 36 or sz > 44:
            s23_font_bad = True
    deductions.append((-1, "第23页标题字号不符(36-44磅)", s23_font_bad))

    # -5: 第3-23页中任意一页存在宽度超过页面宽度90%、高度超过页面高度90%的图片对象,
    #     且该页主要中文标题不可单独选中编辑
    slide_w_emu_d = prs.slide_width
    slide_h_emu_d = prs.slide_height
    threshold_w_d = slide_w_emu_d * 0.9
    threshold_h_d = slide_h_emu_d * 0.9
    large_img_bad = False
    for i in range(2, 23):  # 第3页~第23页,0-indexed 2..22
        slide = slides[i]
        # 找出页面上是否存在尺寸超过 90% 宽 且 90% 高的图片
        oversized_imgs = []
        for shape in slide.shapes:
            if shape.shape_type != 13:  # 仅图片
                continue
            if shape.width is None or shape.height is None:
                continue
            if shape.width > threshold_w_d and shape.height > threshold_h_d:
                oversized_imgs.append(shape)
        if not oversized_imgs:
            continue
        # 主要中文标题是否可单独选中编辑:页面上存在独立的、含中文的文本 shape(非超大图片本身)
        has_editable_chinese_title = False
        for shape in slide.shapes:
            if shape in oversized_imgs:
                continue
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            if not text.strip():
                continue
            if any('一' <= c <= '鿿' for c in text):
                has_editable_chinese_title = True
                break
        if not has_editable_chinese_title:
            large_img_bad = True
            break
    deductions.append((-5, "存在超大图片且标题不可编辑", large_img_bad))

    # ============ 汇总输出 ============
    dim2_items = []
    total_positive = 0
    max_score = 0
    for score, desc, passed in scores:
        max_score += score
        dim2_items.append({
            "rule": desc,
            "max_delta": score,
            "delta": score if passed else 0,
            "hit": bool(passed),
            "detail": "",
        })
        if passed:
            total_positive += score

    total_deduction = 0
    for score, desc, triggered in deductions:
        dim2_items.append({
            "rule": desc,
            "max_delta": score,
            "delta": score if triggered else 0,
            "hit": bool(triggered),
            "detail": "",
        })
        if triggered:
            total_deduction += score

    final_score = total_positive + total_deduction
    final_score = max(0, final_score)  # 不低于0

    result["dim2_items"] = dim2_items
    result["total_score"] = final_score
    result["max_score"] = max_score
    return result


if __name__ == "__main__":
    _dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    try:
        _result = evaluate(_dir)
    except Exception as e:
        _result = {
            "id": SCRIPT_ID,
            "file_name": "",
            "status": "error",
            "error": f"{type(e).__name__}: {e}",
            "dim1_pass": False,
            "dim1_reason": "",
            "dim2_items": [],
            "total_score": 0,
            "max_score": 0,
        }
    print(json.dumps(_result, ensure_ascii=False, indent=2))
