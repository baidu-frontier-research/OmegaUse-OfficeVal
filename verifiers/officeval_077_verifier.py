# -*- coding: utf-8 -*-
"""
自动评估脚本：启蒙运动思想解放_第一页动画修改版(4).pptx

评估逻辑：
  维度1（可用与可修改性）—— 一票否决。任一条不满足 => 总分 0，不再检查维度2。
  维度2（完成度）—— 仅在维度1全部通过后评估。
      得分点：必须满足该细则内的【每一个】子条件才加分（正分）。
      扣分点：只要命中细则内的【任意一个】子条件即扣分（负分）。
  最终累计所有命中细则的分数，并打印命中明细与总分。

实现说明（对“不好实现”的点的变通方式）：
  - .pptx 本质是 zip+OOXML。脚本直接解析 ppt/slides/slide1.xml 的
    形状(p:sp)、图片(p:pic) 与动画时间轴(p:timing) 来判定每一次点击触发了什么效果。
  - “初始隐藏”通过形状 cNvPr 上的 hidden="1" 判定。
  - “第N次点击后淡入/淡出某对象”通过解析 mainSeq 中第 N 个 clickEffect 组，
    检查其中是否对目标形状执行了 fade in / fade out 动画来判定。
  - 本模板中“洛克/休谟/亚当·斯密”行以文本形式呈现“各自主张”（无独立肖像图片），
    故对这些点以“被动画的‘各自主张’文本对象”作为评分主体，符合动画意图。
"""

import sys
import os
import re
import json
import zipfile

# 脚本编号（与文件名 officeval_077_verifier.py 对应，用于结构化返回中的 id 字段）
SCRIPT_ID = "077"

NS_RE = {
    "sp": re.compile(r"<p:sp>(.*?)</p:sp>", re.S),
    "pic": re.compile(r"<p:pic>(.*?)</p:pic>", re.S),
    "cnvpr": re.compile(r'<p:cNvPr id="(\d+)" name="([^"]*)"([^>]*)>'),
    "txt": re.compile(r"<a:t>(.*?)</a:t>", re.S),
    "off": re.compile(r'<a:off x="(-?\d+)" y="(-?\d+)"'),
    "ext": re.compile(r'<a:ext cx="(\d+)" cy="(\d+)"'),
}


# --------------------------------------------------------------------------
# 解析 .pptx
# --------------------------------------------------------------------------
class Shape:
    def __init__(self, sid, name, text, hidden, off, ext, is_pic):
        self.sid = sid
        self.name = name
        self.text = text
        self.hidden = hidden
        self.off = off          # (x, y) EMU 或 None
        self.ext = ext          # (cx, cy) EMU 或 None
        self.is_pic = is_pic


class Slide:
    def __init__(self, xml):
        self.xml = xml
        self.shapes = {}        # sid -> Shape
        self.click_groups = []  # list[ list[effect dict] ]  按点击次序
        self.click_start_indefinite = []  # 每个点击组的开始条件是否为 indefinite
        self.has_click_advance = False    # timing 是否含 onClick/onNext 推进条件
        self._parse_shapes()
        self._parse_timing()

    def _parse_one(self, blk, is_pic):
        m = NS_RE["cnvpr"].search(blk)
        if not m:
            return None
        sid, name, extra = m.group(1), m.group(2), m.group(3)
        hidden = 'hidden="1"' in extra
        text = "".join(NS_RE["txt"].findall(blk)).strip()
        om = NS_RE["off"].search(blk)
        em = NS_RE["ext"].search(blk)
        off = (int(om.group(1)), int(om.group(2))) if om else None
        ext = (int(em.group(1)), int(em.group(2))) if em else None
        return Shape(sid, name, text, hidden, off, ext, is_pic)

    def _parse_shapes(self):
        for blk in NS_RE["sp"].findall(self.xml):
            s = self._parse_one(blk, False)
            if s:
                self.shapes[s.sid] = s
        for blk in NS_RE["pic"].findall(self.xml):
            s = self._parse_one(blk, True)
            if s:
                self.shapes[s.sid] = s

    def _parse_timing(self):
        i = self.xml.find("<p:timing>")
        if i < 0:
            return
        t = self.xml[i:]

        # 判定整个交互序列(mainSeq)是否具备“鼠标单击可推进”的触发链：
        #   - clickEffect 的开始条件若为 delay="indefinite"，须由序列层的
        #     onClick（鼠标单击）事件来推进；
        #   - 仅有 onNext/onPrev（需方向键/特定前进操作）不算“单击触发动画”；
        #   - 若 timing 中没有 onClick 推进条件，则 delay="indefinite" 的 clickEffect
        #     无法被单击触发播放(动画失效)。
        self.has_click_advance = bool(re.search(r'evt="onClick"', t))

        # 主序列由若干 nodeType="clickEffect" 的 par 组成，按出现顺序即为点击顺序。
        positions = [m.start() for m in re.finditer(r'nodeType="clickEffect"', t)]
        for ci, p in enumerate(positions):
            end = positions[ci + 1] if ci + 1 < len(positions) else len(t)
            seg = t[p:end]
            # 该 clickEffect 自身的开始条件（取紧随其后的第一个 stCondLst）
            cond_m = re.search(r"<p:stCondLst>(.*?)</p:stCondLst>", seg, re.S)
            cond_txt = cond_m.group(1) if cond_m else ""
            start_indefinite = 'delay="indefinite"' in cond_txt
            self.click_start_indefinite.append(start_indefinite)
            effects = []
            for em in re.finditer(r"<p:(animEffect|set|anim)\b(.*?)</p:\1>", seg, re.S):
                kind, body = em.group(1), em.group(2)
                spid = re.search(r'spid="(\d+)"', body)
                if not spid:
                    continue
                trans = re.search(r'transition="([^"]*)"', body)
                filt = re.search(r'filter="([^"]*)"', body)
                attr = re.search(r"<p:attrName>([^<]*)</p:attrName>", body)
                toval = re.search(r'<p:strVal val="([^"]*)"', body)
                effects.append({
                    "kind": kind,
                    "spid": spid.group(1),
                    "transition": trans.group(1) if trans else "",
                    "filter": filt.group(1) if filt else "",
                    "attr": attr.group(1) if attr else "",
                    "to": toval.group(1) if toval else "",
                })
            self.click_groups.append(effects)

    # 帮助方法 -------------------------------------------------------------
    def click(self, n):
        """返回第 n 次点击(1-based)的效果列表，越界返回 []"""
        if 1 <= n <= len(self.click_groups):
            return self.click_groups[n - 1]
        return []

    def fade_in_ids(self, n):
        return {e["spid"] for e in self.click(n)
                if e["kind"] == "animEffect" and e["transition"] == "in"
                and e["filter"] == "fade"}

    def fade_out_ids(self, n):
        return {e["spid"] for e in self.click(n)
                if e["kind"] == "animEffect" and e["transition"] == "out"
                and e["filter"] == "fade"}

    def find_by_text(self, keyword):
        """返回文本包含 keyword 的所有形状 sid"""
        return [s.sid for s in self.shapes.values() if keyword in s.text]

    def find_by_exact_text(self, text):
        """返回文本【完全等于】text 的所有形状 sid（用于精确定位表头等）"""
        return [s.sid for s in self.shapes.values() if s.text == text]


# --------------------------------------------------------------------------
# 维度1：可用与可修改性（一票否决）
# --------------------------------------------------------------------------
def evaluate_dimension1(pptx_path: str, slides: "list[Slide]"):
    """返回 (ok: bool, details: list[(passed, msg)])"""
    details = []

    # 1.1 文件为 .pptx 且能作为有效的 OOXML 包打开
    ext_ok = pptx_path.lower().endswith(".pptx")
    pkg_ok = False
    open_msg = ""
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        _ = len(prs.slides)
        pkg_ok = True
    except Exception as e:  # 打不开
        open_msg = "（python-pptx 打开失败: %s）" % e
        # 退化判定：能否作为 zip 解出 presentation.xml
        try:
            with zipfile.ZipFile(pptx_path) as z:
                pkg_ok = "ppt/presentation.xml" in z.namelist()
                open_msg += " zip 解析: %s" % ("成功" if pkg_ok else "失败")
        except Exception as e2:
            open_msg += " zip 解析异常: %s" % e2
    # 规则简化：仅要求“交付文件为 .pptx 格式，文件可正常打开”，
    # 不再对可编辑性/可放映性做额外校验（此前的 _check_editable_and_playable 逻辑已弃用）。
    c1 = ext_ok and pkg_ok
    extra_msg = "" if c1 else open_msg
    details.append((c1, "1.1 格式为.pptx且可正常打开 %s" % extra_msg))

    if not slides:
        details.append((False, "1.2 无法读取幻灯片内容"))
        return False, details

    ok = c1
    return ok, details


# --------------------------------------------------------------------------
# 维度2：完成度
# --------------------------------------------------------------------------
def evaluate_dimension2(slides):
    """返回 (score, hits: list[(score, passed, label)])"""
    s1 = slides[0]
    s2 = slides[1] if len(slides) > 1 else None
    SLIDE_W = 12191695  # presentation.xml 中的 sldSz cx
    hits = []

    # 动画失效判定：凡涉及动画的细则，若其依赖的点击动画无法被触发播放，则该条不通过。
    #   失效条件（满足其一即视为该次点击动画失效）：
    #     - 整个交互序列缺少 onClick/onNext 推进条件（has_click_advance=False），
    #       导致 delay="indefinite" 的 clickEffect 永远无法前进；
    #     - 该点击组不存在（越界）。
    def anim_alive(click_n):
        if click_n < 1 or click_n > len(s1.click_groups):
            return False
        start_indef = (s1.click_start_indefinite[click_n - 1]
                       if click_n - 1 < len(s1.click_start_indefinite) else False)
        # indefinite 开始且无点击/翻页推进 → 无法触发播放
        if start_indef and not s1.has_click_advance:
            return False
        return True

    # ---- 得分点 (+3 each) ----
    # 1) +3 第1页洛克图片对象。细则逐点拆解（全部满足才 +3）：
    #    (a) 洛克图片初始隐藏；
    #    (b) 第一次点击后该图片淡入显示；
    #    (c) 洛克图片下方“各自主张”含 自然权利/财产权/有限政府/政府同意 四项；
    #    (d) 该主张内容第一次点击后淡入显示；
    #    (e) 主张与图片同步或紧随其后（同在第1次点击触发）。
    def eval_locke():
        sub = []
        in1 = s1.fade_in_ids(1)  # 第1次点击淡入的对象集合

        # 定位“洛克图片”：洛克行(代表人物=洛克)下方的图片对象。
        locke_row = s1.find_by_text("洛克")
        locke_pic = None
        for sid, sh in s1.shapes.items():
            if not sh.is_pic or not sh.off:
                continue
            # 同列(x接近)且位于洛克名称下方(y更大)
            for r in locke_row:
                rsh = s1.shapes[r]
                if rsh.off and abs(sh.off[0] - rsh.off[0]) < 900000 and sh.off[1] >= rsh.off[1]:
                    locke_pic = sid
                    break
            if locke_pic:
                break

        # (a) 图片初始隐藏
        a = locke_pic is not None and s1.shapes[locke_pic].hidden
        sub.append((a, "(a)洛克图片初始隐藏"))
        # (b) 第一次点击后图片淡入
        b = locke_pic is not None and locke_pic in in1
        sub.append((b, "(b)第1次点击后洛克图片淡入"))

        # (c) 主张文本含全部四项；若同时能定位到洛克图片，则优先选取“位于洛克图片下方
        #     (同列 x 接近，且 y 更大)”的候选，避免命中页面上其他包含相同关键词的说明性文本。
        claim_sid = None
        claim_candidates = [sid for sid, sh in s1.shapes.items()
                            if all(k in sh.text for k in ["自然权利", "财产权", "有限政府", "政府同意"])]
        if locke_pic is not None:
            psh = s1.shapes[locke_pic]
            for sid in claim_candidates:
                csh = s1.shapes[sid]
                if psh.off and csh.off \
                        and abs(csh.off[0] - psh.off[0]) < 2500000 \
                        and csh.off[1] >= psh.off[1]:
                    claim_sid = sid
                    break
        if claim_sid is None and claim_candidates:
            # 未能定位到图片、或没有位置合适的候选，仍记录首个含四项关键词的文本
            # 供 (f) 报告“主张不在图片下方”的原因
            claim_sid = claim_candidates[0]
        c = claim_sid is not None
        sub.append((c, "(c)洛克主张含自然权利/财产权/有限政府/政府同意"))
        # (d) 主张第一次点击后淡入
        dd = claim_sid is not None and claim_sid in in1
        sub.append((dd, "(d)第1次点击后洛克主张淡入"))
        # (e) 主张与图片同步或紧随（二者均由第1次点击触发）
        e = b and dd
        sub.append((e, "(e)主张与图片同步/紧随(同在第1次点击)"))
        # (f) 主张位于洛克图片下方（同列 x 接近，且 y 更大）——rubric 明确要求“图片下方”
        f = False
        if locke_pic is not None and claim_sid is not None:
            psh, csh = s1.shapes[locke_pic], s1.shapes[claim_sid]
            if psh.off and csh.off:
                f = abs(csh.off[0] - psh.off[0]) < 2500000 and csh.off[1] >= psh.off[1]
        sub.append((f, "(f)主张位于洛克图片下方(同列且y更大)"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_locke()
    hits.append((3, ok and anim_alive(1), "第1次点击 洛克图片对象(图片初始隐藏→淡入 + 主张同步淡入)" + detail))

    # 2) +3 第1页孟德斯鸠材料展示层。细则逐点拆解（全部满足才 +3）：
    #    第二次点击时，当前页(第1页)淡入显示【第二页右侧】孟德斯鸠肖像图片 和 观点提炼内容。
    #    拆点：(a)发生在第2次点击；(b)为淡入(fade in)；(c)淡入对象中含孟德斯鸠肖像图片；
    #          (d)淡入对象中含“观点提炼”内容；(e)所淡入的肖像图片与观点提炼，
    #          与第二页【右侧】(x>页宽一半) 的孟德斯鸠肖像/观点提炼内容一致。
    def eval_montesquieu_layer():
        sub = []
        in2 = s1.fade_in_ids(2)  # 第2次点击淡入的对象（(a)(b)：第2次点击的淡入集合）

        # 收集第二页右侧的孟德斯鸠材料：肖像图片 + 观点提炼正文。
        s2_pic = []          # 右侧图片
        s2_view_bodies = []  # 右侧“观点提炼”正文文本
        s2_has_view_title = False
        s2_has_name = False
        if s2:
            # 先确认右侧确为孟德斯鸠（右侧出现“孟德斯鸠”文本）
            for sh in s2.shapes.values():
                if not sh.off:
                    continue
                if sh.off[0] <= SLIDE_W / 2:
                    continue  # 仅看右侧
                if "孟德斯鸠" in sh.text:
                    s2_has_name = True
                if sh.is_pic:
                    s2_pic.append(sh)
                if sh.text == "观点提炼":
                    s2_has_view_title = True
                # 观点提炼正文：紧随标题、含分权要点的长文本
                if ("分权" in sh.text or "立法" in sh.text) and "观点提炼" not in sh.text:
                    s2_view_bodies.append(sh.text)

        # (a)(b) 第2次点击存在淡入对象
        ab = len(in2) > 0
        sub.append((ab, "(a/b)第2次点击为淡入(fade in)"))

        # (c) 淡入对象中含孟德斯鸠肖像图片
        c = any(s1.shapes.get(i) and s1.shapes[i].is_pic for i in in2)
        sub.append((c, "(c)淡入含孟德斯鸠肖像图片"))

        # (d) 淡入对象中含“观点提炼”内容
        d_title = any(s1.shapes.get(i) and "观点提炼" in s1.shapes[i].text for i in in2)
        d_body = any(s1.shapes.get(i) and ("分权" in s1.shapes[i].text or "立法" in s1.shapes[i].text)
                     for i in in2)
        dd = d_title and d_body
        sub.append((dd, "(d)淡入含观点提炼(标题+正文)"))

        # (e) 与第二页右侧孟德斯鸠肖像/观点提炼内容一致
        s2_ok = bool(s2) and s2_has_name and len(s2_pic) > 0 and s2_has_view_title and len(s2_view_bodies) > 0
        # 第1页淡入的观点提炼正文，应与第二页右侧正文内容一致
        in2_view_texts = [s1.shapes[i].text for i in in2
                          if s1.shapes.get(i) and ("分权" in s1.shapes[i].text or "立法" in s1.shapes[i].text)]
        def norm(x):
            return re.sub(r"[\s•·\.。，,、；;]", "", x)
        body_match = any(norm(a) == norm(b) for a in in2_view_texts for b in s2_view_bodies) \
            or any(norm(a) in norm(b) or norm(b) in norm(a)
                   for a in in2_view_texts for b in s2_view_bodies)
        e = s2_ok and body_match
        sub.append((e, "(e)与第二页右侧孟德斯鸠肖像/观点提炼内容一致"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_montesquieu_layer()
    hits.append((3, ok and anim_alive(2), "第2次点击 孟德斯鸠材料展示层(淡入第二页右侧肖像图片+观点提炼)" + detail))

    # 3) +3 第1页孟德斯鸠材料退出。细则逐点拆解（全部满足才 +3）：
    #    第三次点击后，孟德斯鸠肖像 和 观点提炼内容 淡出消失。
    #    拆点：(a)发生在第3次点击；(b)为淡出(fade out)；
    #          (c)淡出对象中含孟德斯鸠肖像图片；(d)淡出对象中含观点提炼内容。
    def eval_montesquieu_exit():
        sub = []
        out3 = s1.fade_out_ids(3)  # 第3次点击淡出的对象集合

        # (a)(b) 第3次点击存在淡出对象
        ab = len(out3) > 0
        sub.append((ab, "(a/b)第3次点击为淡出(fade out)"))

        # (c) 淡出对象中含孟德斯鸠肖像图片
        #     （限定为第2次点击淡入的那张肖像图片，即孟德斯鸠材料层图片）
        m_pic_ids = {i for i in s1.fade_in_ids(2)
                     if s1.shapes.get(i) and s1.shapes[i].is_pic}
        c = any(i in out3 for i in m_pic_ids) if m_pic_ids \
            else any(s1.shapes.get(i) and s1.shapes[i].is_pic for i in out3)
        sub.append((c, "(c)淡出含孟德斯鸠肖像图片"))

        # (d) 淡出对象中含观点提炼内容（标题+正文）
        d_title = any(s1.shapes.get(i) and "观点提炼" in s1.shapes[i].text for i in out3)
        d_body = any(s1.shapes.get(i) and ("分权" in s1.shapes[i].text or "立法" in s1.shapes[i].text)
                     for i in out3)
        dd = d_title and d_body
        sub.append((dd, "(d)淡出含观点提炼(标题+正文)"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_montesquieu_exit()
    hits.append((3, ok and anim_alive(3), "第3次点击 孟德斯鸠材料退出(肖像+观点提炼淡出消失)" + detail))

    # 4) +3 第1页孟德斯鸠图片与主张。细则逐点拆解（全部满足才 +3）：
    #    第四次点击后，孟德斯鸠图片 及 图片下方“分权制衡、法律理性”等各自主张 淡入显示。
    #    拆点：(a)发生在第4次点击；(b)为淡入(fade in)；
    #          (c)淡入对象含孟德斯鸠图片；
    #          (d)淡入对象含“各自主张”且含 分权制衡 与 法律理性；
    #          (e)主张位于该图片下方(同列、y更大)。
    def eval_montesquieu_pic_claim():
        sub = []
        in4 = s1.fade_in_ids(4)  # 第4次点击淡入集合

        # (a)(b) 第4次点击存在淡入
        ab = len(in4) > 0
        sub.append((ab, "(a/b)第4次点击为淡入(fade in)"))

        # (c) 淡入对象含孟德斯鸠图片
        pic_id = None
        for i in in4:
            sh = s1.shapes.get(i)
            if sh and sh.is_pic:
                pic_id = i
                break
        c = pic_id is not None
        sub.append((c, "(c)淡入含孟德斯鸠图片"))

        # (d) 淡入对象含“各自主张”且含 分权制衡 与 法律理性
        claim_id = None
        for i in in4:
            sh = s1.shapes.get(i)
            if sh and "各自主张" in sh.text and "分权制衡" in sh.text and "法律理性" in sh.text:
                claim_id = i
                break
        dd = claim_id is not None
        sub.append((dd, "(d)淡入含主张(各自主张:分权制衡/法律理性)"))

        # (e) 主张位于图片下方（同列 x 接近，且 y 更大）
        e = False
        if pic_id and claim_id:
            psh, csh = s1.shapes[pic_id], s1.shapes[claim_id]
            if psh.off and csh.off:
                e = abs(csh.off[0] - psh.off[0]) < 2500000 and csh.off[1] >= psh.off[1]
        sub.append((e, "(e)主张位于孟德斯鸠图片下方"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_montesquieu_pic_claim()
    hits.append((3, ok and anim_alive(4), "第4次点击 孟德斯鸠图片与主张(图片+下方主张淡入)" + detail))

    # 5) +3 第1页伏尔泰材料展示层。细则逐点拆解（全部满足才 +3）：
    #    第五次点击后，当前页(第1页)淡入显示【第二页左侧】伏尔泰肖像图片 和 观点提炼内容。
    #    拆点：(a)发生在第5次点击；(b)为淡入(fade in)；(c)淡入对象中含伏尔泰肖像图片；
    #          (d)淡入对象中含“观点提炼”内容；(e)所淡入的肖像/观点提炼，
    #          与第二页【左侧】(x<=页宽一半) 的伏尔泰肖像/观点提炼内容一致。
    def eval_voltaire_layer():
        sub = []
        in5 = s1.fade_in_ids(5)

        # 收集第二页左侧伏尔泰材料：肖像图片 + 观点提炼正文 + 名称 + 标题
        v_pic, v_view_bodies = [], []
        v_has_view_title, v_has_name = False, False
        if s2:
            for sh in s2.shapes.values():
                if not sh.off or sh.off[0] > SLIDE_W / 2:
                    continue  # 仅看左侧
                if "伏尔泰" in sh.text:
                    v_has_name = True
                if sh.is_pic:
                    v_pic.append(sh)
                if sh.text == "观点提炼":
                    v_has_view_title = True
                if ("思想" in sh.text or "言论自由" in sh.text) and "观点提炼" not in sh.text and len(sh.text) > 10:
                    v_view_bodies.append(sh.text)

        # (a)(b) 第5次点击存在淡入
        ab = len(in5) > 0
        sub.append((ab, "(a/b)第5次点击为淡入(fade in)"))
        # (c) 淡入含伏尔泰肖像图片
        c = any(s1.shapes.get(i) and s1.shapes[i].is_pic for i in in5)
        sub.append((c, "(c)淡入含伏尔泰肖像图片"))
        # (d) 淡入含观点提炼（标题+正文）
        d_title = any(s1.shapes.get(i) and "观点提炼" in s1.shapes[i].text for i in in5)
        d_body = any(s1.shapes.get(i) and ("言论自由" in s1.shapes[i].text or "思想" in s1.shapes[i].text)
                     and len(s1.shapes[i].text) > 10 for i in in5)
        dd = d_title and d_body
        sub.append((dd, "(d)淡入含观点提炼(标题+正文)"))
        # (e) 与第二页左侧伏尔泰肖像/观点提炼内容一致
        s2_ok = bool(s2) and v_has_name and len(v_pic) > 0 and v_has_view_title and len(v_view_bodies) > 0
        in5_view_texts = [s1.shapes[i].text for i in in5
                          if s1.shapes.get(i) and ("言论自由" in s1.shapes[i].text or "思想" in s1.shapes[i].text)
                          and len(s1.shapes[i].text) > 10]

        def norm(x):
            return re.sub(r"[\s•·\.。，,、；;]", "", x)
        body_match = any(norm(a) == norm(b) or norm(a) in norm(b) or norm(b) in norm(a)
                         for a in in5_view_texts for b in v_view_bodies)
        e = s2_ok and body_match
        sub.append((e, "(e)与第二页左侧伏尔泰肖像/观点提炼内容一致"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_voltaire_layer()
    hits.append((3, ok and anim_alive(5), "第5次点击 伏尔泰材料展示层(淡入第二页左侧肖像图片+观点提炼)" + detail))

    # 6) +3 第1页伏尔泰材料退出。细则逐点拆解（全部满足才 +3）：
    #    第六次点击后，伏尔泰肖像图片 和 观点提炼内容材料 淡出消失。
    #    拆点：(a)发生在第6次点击；(b)为淡出(fade out)；
    #          (c)淡出对象中含伏尔泰肖像图片；(d)淡出对象中含观点提炼内容。
    def eval_voltaire_exit():
        sub = []
        out6 = s1.fade_out_ids(6)

        ab = len(out6) > 0
        sub.append((ab, "(a/b)第6次点击为淡出(fade out)"))

        # (c) 含伏尔泰肖像图片（限定为第5次点击淡入的那张材料层图片，确保进退对应同一张）
        v_pic_ids = {i for i in s1.fade_in_ids(5)
                     if s1.shapes.get(i) and s1.shapes[i].is_pic}
        c = any(i in out6 for i in v_pic_ids) if v_pic_ids \
            else any(s1.shapes.get(i) and s1.shapes[i].is_pic for i in out6)
        sub.append((c, "(c)淡出含伏尔泰肖像图片"))

        # (d) 含观点提炼内容（标题+正文）
        d_title = any(s1.shapes.get(i) and "观点提炼" in s1.shapes[i].text for i in out6)
        d_body = any(s1.shapes.get(i) and ("言论自由" in s1.shapes[i].text or "思想" in s1.shapes[i].text)
                     and len(s1.shapes[i].text) > 10 for i in out6)
        dd = d_title and d_body
        sub.append((dd, "(d)淡出含观点提炼(标题+正文)"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_voltaire_exit()
    hits.append((3, ok and anim_alive(6), "第6次点击 伏尔泰材料退出(肖像图片+观点提炼淡出消失)" + detail))

    # 7) +3 第1页伏尔泰图片与主张。细则逐点拆解（全部满足才 +3）：
    #    第七次点击后，伏尔泰图片 及 图片下方“思想自由、批判专制”等各自主张 淡入显示。
    #    拆点：(a)发生在第7次点击；(b)为淡入(fade in)；
    #          (c)淡入对象含伏尔泰图片；
    #          (d)淡入对象含“各自主张”且含 思想自由 与 批判专制；
    #          (e)主张位于该图片下方(同列、y更大)。
    def eval_voltaire_pic_claim():
        sub = []
        in7 = s1.fade_in_ids(7)

        ab = len(in7) > 0
        sub.append((ab, "(a/b)第7次点击为淡入(fade in)"))

        # (c) 淡入对象含伏尔泰图片
        pic_id = None
        for i in in7:
            sh = s1.shapes.get(i)
            if sh and sh.is_pic:
                pic_id = i
                break
        c = pic_id is not None
        sub.append((c, "(c)淡入含伏尔泰图片"))

        # (d) 淡入对象含“各自主张”且含 思想自由 与 批判专制
        claim_id = None
        for i in in7:
            sh = s1.shapes.get(i)
            if sh and "各自主张" in sh.text and "思想自由" in sh.text and "批判专制" in sh.text:
                claim_id = i
                break
        dd = claim_id is not None
        sub.append((dd, "(d)淡入含主张(各自主张:思想自由/批判专制)"))

        # (e) 主张位于伏尔泰图片下方（同列 x 接近，且 y 更大）
        e = False
        if pic_id and claim_id:
            psh, csh = s1.shapes[pic_id], s1.shapes[claim_id]
            if psh.off and csh.off:
                e = abs(csh.off[0] - psh.off[0]) < 2500000 and csh.off[1] >= psh.off[1]
        sub.append((e, "(e)主张位于伏尔泰图片下方"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_voltaire_pic_claim()
    hits.append((3, ok and anim_alive(7), "第7次点击 伏尔泰图片与主张(图片+下方主张淡入)" + detail))

    # 8) +3 第1页卢梭材料展示层。细则逐点拆解（全部满足才 +3）：
    #    第八次点击后，当前页(第1页)淡入显示【第三页左侧】卢梭肖像图片 和 观点提炼内容。
    #    拆点：(a)发生在第8次点击；(b)为淡入(fade in)；(c)淡入对象含卢梭肖像图片；
    #          (d)淡入对象含“观点提炼”内容；(e)所淡入肖像/观点提炼，
    #          与第三页【左侧】(x<=页宽一半) 的卢梭肖像/观点提炼内容一致。
    def eval_rousseau_layer():
        sub = []
        in8 = s1.fade_in_ids(8)
        s3 = slides[2] if len(slides) > 2 else None

        r_pic, r_view_bodies = [], []
        r_has_view_title, r_has_name = False, False
        if s3:
            for sh in s3.shapes.values():
                if not sh.off or sh.off[0] > SLIDE_W / 2:
                    continue  # 仅看左侧
                if "卢梭" in sh.text:
                    r_has_name = True
                if sh.is_pic:
                    r_pic.append(sh)
                if sh.text == "观点提炼":
                    r_has_view_title = True
                if ("社会契约" in sh.text or "共同体" in sh.text) and "观点提炼" not in sh.text and len(sh.text) > 10:
                    r_view_bodies.append(sh.text)

        ab = len(in8) > 0
        sub.append((ab, "(a/b)第8次点击为淡入(fade in)"))
        c = any(s1.shapes.get(i) and s1.shapes[i].is_pic for i in in8)
        sub.append((c, "(c)淡入含卢梭肖像图片"))
        d_title = any(s1.shapes.get(i) and "观点提炼" in s1.shapes[i].text for i in in8)
        d_body = any(s1.shapes.get(i) and ("社会契约" in s1.shapes[i].text or "共同体" in s1.shapes[i].text)
                     and len(s1.shapes[i].text) > 10 for i in in8)
        dd = d_title and d_body
        sub.append((dd, "(d)淡入含观点提炼(标题+正文)"))
        s3_ok = bool(s3) and r_has_name and len(r_pic) > 0 and r_has_view_title and len(r_view_bodies) > 0
        in8_view_texts = [s1.shapes[i].text for i in in8
                          if s1.shapes.get(i) and ("社会契约" in s1.shapes[i].text or "共同体" in s1.shapes[i].text)
                          and len(s1.shapes[i].text) > 10]

        def norm(x):
            return re.sub(r"[\s•·\.。，,、；;]", "", x)
        body_match = any(norm(a) == norm(b) or norm(a) in norm(b) or norm(b) in norm(a)
                         for a in in8_view_texts for b in r_view_bodies)
        e = s3_ok and body_match
        sub.append((e, "(e)与第三页左侧卢梭肖像/观点提炼内容一致"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_rousseau_layer()
    hits.append((3, ok and anim_alive(8), "第8次点击 卢梭材料展示层(淡入第三页左侧肖像图片+观点提炼)" + detail))

    # 9) +3 第1页卢梭材料退出。细则逐点拆解（全部满足才 +3）：
    #    第九次点击后，卢梭肖像图片 和 观点提炼内容 淡出消失。
    #    拆点：(a)发生在第9次点击；(b)为淡出(fade out)；
    #          (c)淡出对象中含卢梭肖像图片；(d)淡出对象中含观点提炼内容。
    def eval_rousseau_exit():
        sub = []
        out9 = s1.fade_out_ids(9)

        ab = len(out9) > 0
        sub.append((ab, "(a/b)第9次点击为淡出(fade out)"))

        # (c) 含卢梭肖像图片（限定为第8次点击淡入的那张材料层图片，确保进退对应同一张）
        r_pic_ids = {i for i in s1.fade_in_ids(8)
                     if s1.shapes.get(i) and s1.shapes[i].is_pic}
        c = any(i in out9 for i in r_pic_ids) if r_pic_ids \
            else any(s1.shapes.get(i) and s1.shapes[i].is_pic for i in out9)
        sub.append((c, "(c)淡出含卢梭肖像图片"))

        # (d) 含观点提炼内容（标题+正文）
        d_title = any(s1.shapes.get(i) and "观点提炼" in s1.shapes[i].text for i in out9)
        d_body = any(s1.shapes.get(i) and ("社会契约" in s1.shapes[i].text or "共同体" in s1.shapes[i].text)
                     and len(s1.shapes[i].text) > 10 for i in out9)
        dd = d_title and d_body
        sub.append((dd, "(d)淡出含观点提炼(标题+正文)"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_rousseau_exit()
    hits.append((3, ok and anim_alive(9), "第9次点击 卢梭材料退出(肖像图片+观点提炼淡出消失)" + detail))

    # 10) +3 第1页卢梭图片与主张。细则逐点拆解（全部满足才 +3）：
    #    第十次点击后，卢梭图片 及 图片下方“社会契约、主权在民”等各自主张 淡入显示。
    #    拆点：(a)发生在第10次点击；(b)为淡入(fade in)；
    #          (c)淡入对象含卢梭图片；
    #          (d)淡入对象含“各自主张”且含 社会契约 与 主权在民；
    #          (e)主张位于该图片下方(同列、y更大)。
    def eval_rousseau_pic_claim():
        sub = []
        in10 = s1.fade_in_ids(10)

        ab = len(in10) > 0
        sub.append((ab, "(a/b)第10次点击为淡入(fade in)"))

        pic_id = None
        for i in in10:
            sh = s1.shapes.get(i)
            if sh and sh.is_pic:
                pic_id = i
                break
        c = pic_id is not None
        sub.append((c, "(c)淡入含卢梭图片"))

        claim_id = None
        for i in in10:
            sh = s1.shapes.get(i)
            if sh and "各自主张" in sh.text and "社会契约" in sh.text and "主权在民" in sh.text:
                claim_id = i
                break
        dd = claim_id is not None
        sub.append((dd, "(d)淡入含主张(各自主张:社会契约/主权在民)"))

        e = False
        if pic_id and claim_id:
            psh, csh = s1.shapes[pic_id], s1.shapes[claim_id]
            if psh.off and csh.off:
                e = abs(csh.off[0] - psh.off[0]) < 2500000 and csh.off[1] >= psh.off[1]
        sub.append((e, "(e)主张位于卢梭图片下方"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_rousseau_pic_claim()
    hits.append((3, ok and anim_alive(10), "第10次点击 卢梭图片与主张(图片+下方主张淡入)" + detail))

    # 11) +3 第1页休谟主张文本。细则逐点拆解（全部满足才 +3）：
    #     第十一次点击后，休谟图片下方“经验主义、理性边界”等各自主张内容
    #     与休谟图片同步或紧随其后淡入显示。
    #     拆点：(a)发生在第11次点击；(b)为淡入(fade in)；
    #           (c)淡入对象含“各自主张”且含 经验主义 与 理性(边界)；
    #           (d)存在休谟图片，且该主张位于图片下方(同列、y更大)；
    #           (e)主张与休谟图片同步或紧随（图片亦在第11次点击淡入）。
    def eval_hume_claim():
        sub = []
        in11 = s1.fade_in_ids(11)

        ab = len(in11) > 0
        sub.append((ab, "(a/b)第11次点击为淡入(fade in)"))

        # (c) 主张含 经验主义 与 理性边界（rubric 要求“理性边界”完整词，
        #     不能仅命中“理性”避免把“理性主义”等文本误判为合格）
        claim_id = None
        for i in in11:
            sh = s1.shapes.get(i)
            if sh and "各自主张" in sh.text and "经验主义" in sh.text and "理性边界" in sh.text:
                claim_id = i
                break
        c = claim_id is not None
        sub.append((c, "(c)淡入含主张(各自主张:经验主义/理性边界)"))

        # 定位休谟图片：休谟名称行下方/同列的图片对象
        hume_rows = [sid for sid, sh in s1.shapes.items() if sh.text == "休谟"]
        hume_pic = None
        for sid, sh in s1.shapes.items():
            if not sh.is_pic or not sh.off:
                continue
            for r in hume_rows:
                rsh = s1.shapes[r]
                if rsh.off and abs(sh.off[0] - rsh.off[0]) < 900000 and sh.off[1] >= rsh.off[1]:
                    hume_pic = sid
                    break
            if hume_pic:
                break

        # (d) 存在休谟图片，且主张位于其下方
        dd = False
        if hume_pic and claim_id:
            psh, csh = s1.shapes[hume_pic], s1.shapes[claim_id]
            if psh.off and csh.off:
                dd = abs(csh.off[0] - psh.off[0]) < 2500000 and csh.off[1] >= psh.off[1]
        sub.append((dd, "(d)休谟图片存在且主张位于图片下方"))

        # (e) 主张与休谟图片同步/紧随（图片亦在第11次点击淡入）
        e = hume_pic is not None and hume_pic in in11 and c
        sub.append((e, "(e)主张与休谟图片同步/紧随(同在第11次点击淡入)"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_hume_claim()
    hits.append((3, ok and anim_alive(11), "第11次点击 休谟主张文本(休谟图片下方主张与图片同步/紧随淡入)" + detail))

    # 12) +3 第1页亚当·斯密图片对象。细则逐点拆解（全部满足才 +3）：
    #     第十二次点击后，亚当·斯密图片 及 下方“分工市场、自由竞争”等各自主张内容 淡入显示。
    #     拆点：(a)发生在第12次点击；(b)为淡入(fade in)；
    #           (c)淡入对象含亚当·斯密图片(需绑定为亚当·斯密名称行下方的图片)；
    #           (d)淡入对象含“各自主张”且含完整词 分工市场 与 自由竞争；
    #           (e)主张位于该图片下方(同列、y更大)。
    def eval_smith_pic_claim():
        sub = []
        in12 = s1.fade_in_ids(12)

        ab = len(in12) > 0
        sub.append((ab, "(a/b)第12次点击为淡入(fade in)"))

        # 先定位“亚当·斯密图片”：亚当·斯密名称行下方/同列的图片对象。
        #   名称文本存在两种常见写法(中点/圆点)，两种都容纳。
        smith_rows = [sid for sid, sh in s1.shapes.items()
                      if "亚当·斯密" in sh.text or "亚当・斯密" in sh.text or "亚当斯密" in sh.text]
        smith_pic_all = None
        for sid, sh in s1.shapes.items():
            if not sh.is_pic or not sh.off:
                continue
            for r in smith_rows:
                rsh = s1.shapes[r]
                if rsh.off and abs(sh.off[0] - rsh.off[0]) < 900000 and sh.off[1] >= rsh.off[1]:
                    smith_pic_all = sid
                    break
            if smith_pic_all:
                break

        # (c) 淡入对象含亚当·斯密图片：淡入集合中必须包含定位到的亚当·斯密图片
        pic_id = smith_pic_all if (smith_pic_all is not None and smith_pic_all in in12) else None
        c = pic_id is not None
        sub.append((c, "(c)淡入含亚当·斯密图片(已绑定亚当·斯密名称行下方图片)"))

        # (d) 淡入对象含“各自主张”且含完整词 分工市场 与 自由竞争
        #     rubric 要求完整词“分工市场”，不能仅命中“分工”避免漏掉“市场”二字。
        claim_id = None
        for i in in12:
            sh = s1.shapes.get(i)
            if sh and "各自主张" in sh.text and "分工市场" in sh.text and "自由竞争" in sh.text:
                claim_id = i
                break
        dd = claim_id is not None
        sub.append((dd, "(d)淡入含主张(各自主张:分工市场/自由竞争)"))

        # (e) 主张位于亚当·斯密图片下方
        e = False
        if pic_id and claim_id:
            psh, csh = s1.shapes[pic_id], s1.shapes[claim_id]
            if psh.off and csh.off:
                e = abs(csh.off[0] - psh.off[0]) < 2500000 and csh.off[1] >= psh.off[1]
        sub.append((e, "(e)主张位于亚当·斯密图片下方"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_smith_pic_claim()
    hits.append((3, ok and anim_alive(12), "第12次点击 亚当·斯密图片对象(图片+下方主张淡入)" + detail))

    # 13) +3 第1页康德材料展示层。细则逐点拆解（全部满足才 +3）：
    #     第十三次点击后，当前页(第1页)淡入显示【第四页右侧】康德肖像图片 和 观点提炼内容。
    #     拆点：(a)发生在第13次点击；(b)为淡入(fade in)；(c)淡入对象含康德肖像图片；
    #           (d)淡入对象含“观点提炼”内容；(e)所淡入肖像/观点提炼，
    #           与第四页【右侧】(x>页宽一半) 的康德肖像/观点提炼内容一致。
    def eval_kant_layer():
        sub = []
        in13 = s1.fade_in_ids(13)
        s4 = slides[3] if len(slides) > 3 else None

        k_pic, k_view_bodies = [], []
        k_has_view_title, k_has_name = False, False
        if s4:
            for sh in s4.shapes.values():
                if not sh.off or sh.off[0] <= SLIDE_W / 2:
                    continue  # 仅看右侧
                if "康德" in sh.text:
                    k_has_name = True
                if sh.is_pic:
                    k_pic.append(sh)
                if sh.text == "观点提炼":
                    k_has_view_title = True
                if ("独立思考" in sh.text or "理性" in sh.text) and "观点提炼" not in sh.text and len(sh.text) > 10:
                    k_view_bodies.append(sh.text)

        ab = len(in13) > 0
        sub.append((ab, "(a/b)第13次点击为淡入(fade in)"))
        c = any(s1.shapes.get(i) and s1.shapes[i].is_pic for i in in13)
        sub.append((c, "(c)淡入含康德肖像图片"))
        d_title = any(s1.shapes.get(i) and "观点提炼" in s1.shapes[i].text for i in in13)
        d_body = any(s1.shapes.get(i) and ("独立思考" in s1.shapes[i].text or "理性" in s1.shapes[i].text)
                     and len(s1.shapes[i].text) > 10 for i in in13)
        dd = d_title and d_body
        sub.append((dd, "(d)淡入含观点提炼(标题+正文)"))
        s4_ok = bool(s4) and k_has_name and len(k_pic) > 0 and k_has_view_title and len(k_view_bodies) > 0
        in13_view_texts = [s1.shapes[i].text for i in in13
                           if s1.shapes.get(i) and ("独立思考" in s1.shapes[i].text or "理性" in s1.shapes[i].text)
                           and len(s1.shapes[i].text) > 10]

        def norm(x):
            return re.sub(r"[\s•·\.。，,、；;]", "", x)
        body_match = any(norm(a) == norm(b) or norm(a) in norm(b) or norm(b) in norm(a)
                         for a in in13_view_texts for b in k_view_bodies)
        e = s4_ok and body_match
        sub.append((e, "(e)与第四页右侧康德肖像/观点提炼内容一致"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_kant_layer()
    hits.append((3, ok and anim_alive(13), "第13次点击 康德材料展示层(淡入第四页右侧肖像图片+观点提炼)" + detail))

    # 14) +3 第1页康德材料退出。细则逐点拆解（全部满足才 +3）：
    #     第十四次点击，第四页康德肖像图片 和 观点提炼内容 淡出消失，不遮挡后续康德主张内容。
    #     拆点：(a)发生在第14次点击；(b)为淡出(fade out)；
    #           (c)淡出对象含康德肖像图片；(d)淡出对象含观点提炼内容；
    #           (e)不遮挡后续康德主张：即覆盖在康德主张(id:独立思考/自由自律)区域上方的
    #              康德材料层对象，在本次点击全部淡出（淡出后主张区域不再被材料层遮挡）。
    def eval_kant_exit():
        sub = []
        out14 = s1.fade_out_ids(14)

        ab = len(out14) > 0
        sub.append((ab, "(a/b)第14次点击为淡出(fade out)"))

        # (c) 含康德肖像图片（限定为第13次点击淡入的那张材料层图片）
        k_pic_ids = {i for i in s1.fade_in_ids(13)
                     if s1.shapes.get(i) and s1.shapes[i].is_pic}
        c = any(i in out14 for i in k_pic_ids) if k_pic_ids \
            else any(s1.shapes.get(i) and s1.shapes[i].is_pic for i in out14)
        sub.append((c, "(c)淡出含康德肖像图片"))

        # (d) 含观点提炼内容（标题+正文）
        d_title = any(s1.shapes.get(i) and "观点提炼" in s1.shapes[i].text for i in out14)
        d_body = any(s1.shapes.get(i) and ("独立思考" in s1.shapes[i].text or "理性" in s1.shapes[i].text)
                     and len(s1.shapes[i].text) > 10 for i in out14)
        dd = d_title and d_body
        sub.append((dd, "(d)淡出含观点提炼(标题+正文)"))

        # (e) 不遮挡后续康德主张：找到康德主张文本(各自主张:独立思考/自由自律)的矩形区域，
        #     凡是第13次点击淡入(=材料层)且其矩形与主张区域重叠的对象，本次必须全部淡出。
        def rect(sh):
            if not sh.off or not sh.ext:
                return None
            x, y = sh.off
            cx, cy = sh.ext
            return (x, y, x + cx, y + cy)

        def overlap(a, b):
            return a and b and not (a[2] <= b[0] or b[2] <= a[0] or a[3] <= b[1] or b[3] <= a[1])

        claim_sh = None
        for sid, sh in s1.shapes.items():
            if "各自主张" in sh.text and "独立思考" in sh.text and "自由自律" in sh.text:
                claim_sh = sh
                break
        claim_rect = rect(claim_sh) if claim_sh else None

        layer_ids = set(s1.fade_in_ids(13))  # 第13次淡入的康德材料层对象集合
        occluders = [i for i in layer_ids
                     if claim_rect and overlap(rect(s1.shapes.get(i)), claim_rect)]
        # 不遮挡条件：所有遮挡主张区域的材料层对象都在第14次点击被淡出
        e = claim_rect is not None and all(i in out14 for i in occluders)
        sub.append((e, "(e)遮挡康德主张区域的材料层对象均已淡出(不遮挡后续主张)"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_kant_exit()
    hits.append((3, ok and anim_alive(14), "第14次点击 康德材料退出(肖像+观点提炼淡出消失,不遮挡后续主张)" + detail))

    # 15) +3 第1页康德图片与主张。细则逐点拆解（全部满足才 +3）：
    #     第十五次点击后，第一页康德图片 及 图片下方“独立思考、自由自律”等各自主张 淡入显示。
    #     拆点：(a)发生在第15次点击；(b)为淡入(fade in)；
    #           (c)淡入对象含康德图片；
    #           (d)淡入对象含“各自主张”且含 独立思考 与 自由自律；
    #           (e)主张位于该图片下方(同列、y更大)。
    def eval_kant_pic_claim():
        sub = []
        in15 = s1.fade_in_ids(15)

        ab = len(in15) > 0
        sub.append((ab, "(a/b)第15次点击为淡入(fade in)"))

        pic_id = None
        for i in in15:
            sh = s1.shapes.get(i)
            if sh and sh.is_pic:
                pic_id = i
                break
        c = pic_id is not None
        sub.append((c, "(c)淡入含康德图片"))

        claim_id = None
        for i in in15:
            sh = s1.shapes.get(i)
            if sh and "各自主张" in sh.text and "独立思考" in sh.text and "自由自律" in sh.text:
                claim_id = i
                break
        dd = claim_id is not None
        sub.append((dd, "(d)淡入含主张(各自主张:独立思考/自由自律)"))

        e = False
        if pic_id and claim_id:
            psh, csh = s1.shapes[pic_id], s1.shapes[claim_id]
            if psh.off and csh.off:
                e = abs(csh.off[0] - psh.off[0]) < 2500000 and csh.off[1] >= psh.off[1]
        sub.append((e, "(e)主张位于康德图片下方"))

        ok = all(p for p, _ in sub)
        miss = [lab for p, lab in sub if not p]
        detail = "" if ok else "（未满足: %s）" % "; ".join(miss)
        return ok, detail

    ok, detail = eval_kant_pic_claim()
    hits.append((3, ok and anim_alive(15), "第15次点击 康德图片与主张(图片+下方主张淡入)" + detail))

    # ---- 扣分点 ----
    # -5：第1页在放映时未出现思想脉络板块内容。
    #     判定（内容口径，不强制标题）：只要放映时(初始可见状态)第1页出现该板块的
    #       主要内容之一，即视为板块已出现；不再单独要求“思想脉络”四字标题存在。
    #       命中(扣 -5)条件：板块表头(时间/国家/代表人物) 与 7 位人物行(洛克/孟德斯鸠/
    #         伏尔泰/卢梭/休谟/亚当·斯密/康德) 大部分都不可见 —— 具体判定为：
    #         可见表头字段数 + 可见人物行数 < 4 （即板块骨架缺失）。
    #       “放映时可见” = 形状存在 且 初始未隐藏(hidden=False)。
    def visible_contains(text):
        return [sid for sid, sh in s1.shapes.items()
                if text in sh.text and not sh.hidden]

    deduct1 = False
    reasons1 = []

    header_fields = ["时间", "国家", "代表人物"]
    persons_all = ["洛克", "孟德斯鸠", "伏尔泰", "卢梭", "休谟", "亚当·斯密", "康德"]
    visible_headers = [h for h in header_fields if visible_contains(h)]
    visible_persons = [p for p in persons_all if visible_contains(p)]
    has_block_title = len(visible_contains("思想脉络")) > 0
    # 板块内容存在的判定：标题、任一表头字段、或任一人物行可见，均视为板块已出现。
    has_block_content = has_block_title or bool(visible_headers) or bool(visible_persons)
    # 严格骨架校验：可见表头 + 可见人物行合计过少，认定板块内容缺失。
    skeleton_ok = (len(visible_headers) + len(visible_persons)) >= 4

    if not has_block_content or not skeleton_ok:
        deduct1 = True
        if not visible_headers:
            reasons1.append("表头(时间/国家/代表人物)均不可见")
        else:
            missing_h = [h for h in header_fields if h not in visible_headers]
            if missing_h:
                reasons1.append("表头缺失: %s" % "/".join(missing_h))
        missing_p = [p for p in persons_all if p not in visible_persons]
        if missing_p:
            reasons1.append("人物行缺失: %s" % "/".join(missing_p))
        if not reasons1:
            reasons1.append("第1页未出现可识别的思想脉络板块内容")

    hits.append((-5, deduct1, "第1页在放映时未出现思想脉络板块内容"
                 + ("（%s）" % "; ".join(reasons1[:3]) if deduct1 else "")))

    score = sum(sc for sc, ok, _ in hits if ok)
    return score, hits


# --------------------------------------------------------------------------
# 主流程（统一接口：见同级/上级目录《脚本接口差异与统一建议.md》§2）
# --------------------------------------------------------------------------
def load_slides(pptx_path):
    slides = []
    with zipfile.ZipFile(pptx_path) as z:
        names = sorted([n for n in z.namelist()
                        if re.match(r"ppt/slides/slide\d+\.xml$", n)],
                       key=lambda n: int(re.search(r"(\d+)", n.split("/")[-1]).group(1)))
        for n in names:
            slides.append(Slide(z.read(n).decode("utf-8")))
    return slides


def _find_pptx_in_dir(dir_path: str):
    """在给定目录中定位被评估的 .pptx 文档；找不到返回 None。
    忽略 Office 打开时生成的临时文件（形如 ~$xxx.pptx）。"""
    if not os.path.isdir(dir_path):
        return None
    candidates = [n for n in os.listdir(dir_path)
                  if n.lower().endswith(".pptx") and not n.startswith("~$")]
    if not candidates:
        return None
    candidates.sort()
    return os.path.join(dir_path, candidates[0])


def _split_rule_detail(label: str):
    """将 hits 中的 label 拆分为 (rule, detail)。
    detail 取尾部『（未满足: ...）』或『（...）』括号内容；无则为空串。"""
    m = re.search(r"（(?:未满足:\s*)?([^（）]*)）\s*$", label)
    if m:
        return label[:m.start()].rstrip(), m.group(1).strip()
    return label, ""


def evaluate(dir_path: str) -> dict:
    """统一评估入口。

    参数：
        dir_path: 脚本所在目录的路径；脚本自身负责在该目录内定位并打开被评估文档。

    返回：
        结构化字典，字段见《脚本接口差异与统一建议.md》§2.2。
    """
    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": 0,
    }
    try:
        pptx_path = _find_pptx_in_dir(dir_path)
        if not pptx_path or not os.path.exists(pptx_path):
            result["status"] = "error"
            result["error"] = "目录中未找到 .pptx 文件：%s" % dir_path
            return result
        result["file_name"] = os.path.basename(pptx_path)

        try:
            slides = load_slides(pptx_path)
        except Exception as e:
            slides = []
            # 解析失败视为维度1不可用（非脚本自身崩溃），在 dim1 阶段兜底给出原因
            result["dim1_reason"] = "解析幻灯片失败：%s" % e

        d1_ok, d1_details = evaluate_dimension1(pptx_path, slides)
        result["dim1_pass"] = bool(d1_ok)

        if not d1_ok:
            reasons = [msg for passed, msg in d1_details if not passed]
            if reasons:
                # 若前面已因解析失败填过原因，则叠加维度1细则原因
                prev = result["dim1_reason"]
                joined = "; ".join(reasons)
                result["dim1_reason"] = "%s; %s" % (prev, joined) if prev else joined
            # 维度1未通过：dim2_items 为空，total_score=0；仍尝试给出 max_score 便于 Excel 对齐
            try:
                _, sample_hits = evaluate_dimension2(slides)
                result["max_score"] = sum(sc for sc, _, _ in sample_hits if sc > 0)
            except Exception:
                result["max_score"] = 0
            return result

        score, hits = evaluate_dimension2(slides)
        items = []
        max_score = 0
        for sc, ok, label in hits:
            rule, _ = _split_rule_detail(label)
            if sc > 0:
                max_score += sc
            items.append({
                "rule": rule,
                "max_delta": sc,
                "delta": sc if ok else 0,
                "hit": bool(ok),
                "detail": "",
            })
        result["dim2_items"] = items
        result["total_score"] = score
        result["max_score"] = max_score
        return result
    except Exception as e:
        result["status"] = "error"
        result["error"] = "%s: %s" % (type(e).__name__, e)
        return result


if __name__ == "__main__":
    # 本地调试用：默认取脚本自身所在目录；也可显式传入目录路径。
    dbg_dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(dbg_dir), ensure_ascii=False, indent=2))
