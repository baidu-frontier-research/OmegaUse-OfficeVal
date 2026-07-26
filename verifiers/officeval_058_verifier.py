# -*- coding: utf-8 -*-
"""
自动评估脚本：对 PPTX 目录标签作业按"打分细则"进行评分。

评估流程：
  维度1（可用与可修改性）—— 不满足 => 直接 0 分，不再检查维度2。
  维度2（完成度）—— 逐条评分细则检测：
        加分细则：必须满足该细则中的【每一个】点才加分。
        扣分细则：只要满足细则中的【任意一个】点即扣分。
  对外仅暴露 evaluate(dir_path: str) -> dict，返回结构化评估结果。
"""

import os
import sys
import json

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_CM = 360000.0
def to_cm(emu):
    return None if emu is None else emu / EMU_PER_CM

# ---------- 预期文字（四组目录项 × 3 个标签） ----------
EXPECTED = {
    "文化沟通": ["儒家伦理", "叙事逻辑", "文化折扣"],
    "工业体系": ["特效能力", "宣发困境", "产业链断点"],
    "全球竞争": ["好莱坞", "韩国电影", "发行网络"],
    "认知偏差": ["预设立场", "误解误读", "题材受限"],
}
# 目录项文字 -> 大致的 top（cm），用于把标签按行归到对应目录组
# 由文件实测：01=4.88, 02=7.54, 03=10.21, 04=12.88（目录项 Text 顶部）
GROUP_ORDER = ["文化沟通", "工业体系", "全球竞争", "认知偏差"]

# 各加分项对应的"评分细则内容"原文（命中时打印用）
RULE_TEXT = {
    "文化沟通": '第2页"文化沟通"目录项下方从左往右数第1个标签：文字为"儒家伦理"。第2个标签：文字为"叙事逻辑"。第3个标签：文字为"文化折扣"。',
    "工业体系": '第2页"工业体系"目录项下方从左往右数第1个标签：文字为"特效能力"。第2个标签：文字为"宣发困境"。第3个标签：文字为"产业链断点"。',
    "全球竞争": '第2页"全球竞争"目录项下方从左往右数第1个标签：文字为"好莱坞"。第2个标签：文字为"韩国电影"。第3个标签：文字为"发行网络"。',
    "认知偏差": '第2页"认知偏差"目录项下方从左往右数第1个标签：文字为"预设立场"。第2个标签：文字为"误解误读"。第3个标签：文字为"题材受限"。',
    "位置": '第2页标签整体位置：四组标签分别位于对应目录项文字下方的空白区域，不进入右侧图片区域或页面底部页码区域。每个目录项下方恰好放置3个标签，同组标签从左至右排列',
    "尺寸": '第2页全部12个标签尺寸：宽度约1.8至2.5厘米、高度约0.55至0.85厘米，同组及不同组标签大小一致，宽高偏差不超过0.1厘米。水平间距相等，建议保持0.25至0.45厘米，任意两个间距差不超过0.1厘米。',
    "文字": '第2页全部标签文字：字体统一为MiSans，字号统一为10.5磅，文字水平居中、垂直居中。',
    "形状": '第2页全部标签形状：使用胶囊形或高圆角矩形，无明显尖角，边框为无轮廓或与填充色一致。形状使用深绿色填充，文字使用白色，颜色在12个标签中保持一致。',
    "扣-3-尺寸": '任意三个以上标签大小差异超过1厘米。',
    "扣-3-页数": 'PPT页数不是8页。',
    "扣-1": '文件中出现批注、临时说明文字、红色标记、截图边框或多余占位对象。',
}


# =========================================================
# 工具函数
# =========================================================
def normalize(s):
    return (s or "").strip().replace("　", "").replace(" ", "")

def is_rounded_label(sh):
    """判断是否为我们新增的圆角标签。"""
    if sh.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
        return False
    try:
        ast = sh.auto_shape_type
    except Exception:
        return False
    # 胶囊/圆角矩形
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    if ast not in (MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,):
        # 也接受名称包含 Rounded / 文字命中预期的形状
        if "Rounded" not in sh.name:
            return False
    # 必须有文字
    if not sh.has_text_frame:
        return False
    return True


# =========================================================
# 读取文件
# =========================================================
def load(pptx_path):
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return None, f"无法打开文件：{e}"
    return prs, None


# =========================================================
# 维度 1：可用与可修改性
# =========================================================
def check_dimension1(prs):
    """返回 (passed: bool, reasons: list[str], labels: list)

    维度1目前只做“文件能被 python-pptx 正常打开/解析”这一条底线校验——
    load() 已经在打开失败时返回 err，本函数在此基础上仅收集第2页的圆角标签
    供维度2使用，不再对页数、标签数、整页图覆盖、越界排版等做门槛判断。
    """
    reasons = ["✓ 文件为 .pptx 且可正常打开/解析"]
    labels: list[object] = []
    if len(prs.slides) >= 2:
        slide2 = prs.slides[1]
        labels = [sh for sh in slide2.shapes if is_rounded_label(sh)]
    return True, reasons, labels



# =========================================================
# 维度 2 辅助：把 12 个标签归到 4 组（按 top 行聚类，再按 left 排序）
# =========================================================
def group_labels(labels):
    """返回 {组名: [按 left 排序的 label,...]}，按 top 升序的四行映射到 GROUP_ORDER。"""
    if len(labels) != 12:
        # 仍尽力分组
        pass
    # 按 top 聚类成行（容差 1cm）
    items = sorted(labels, key=lambda s: (s.top or 0))
    rows = []
    for l in items:
        placed = False
        for row in rows:
            if abs((l.top or 0) - row[0]) < 1.0 * EMU_PER_CM:
                row[1].append(l); placed = True; break
        if not placed:
            rows.append([(l.top or 0), [l]])
    rows.sort(key=lambda r: r[0])
    groups = {}
    for idx, (_, row) in enumerate(rows):
        if idx < len(GROUP_ORDER):
            gname = GROUP_ORDER[idx]
        else:
            gname = f"额外行{idx+1}"
        groups[gname] = sorted(row, key=lambda s: (s.left or 0))
    return groups


# =========================================================
# 维度 2：完成度评分细则
# =========================================================
def check_dimension2(prs, labels):
    """返回 (total_score, hits: list[(score, desc, detail)])"""
    hits = []
    groups = group_labels(labels)
    slide2 = prs.slides[1]

    # ---------- 加分项 1~4：四组文字正确 ----------
    # 细则原文（以"文化沟通"为例）：
    #   第2页"X"目录项下方从左往右数第1个标签：文字为"…"。
    #   第2个标签：文字为"…"。第3个标签：文字为"…"。
    # 逐点踩：① 位于该目录项"下方"；② 从左往右数第1个=…；③ 第2个=…；④ 第3个=…。
    # 加分细则需"每一个点"都满足才 +3。
    DIR_TOP_TEXT = {  # 各目录项文字（Text 形状）的 top（cm），用于判定"下方"
        "文化沟通": 4.826, "工业体系": 7.493, "全球竞争": 10.16, "认知偏差": 12.827,
    }
    for gname in GROUP_ORDER:
        exp = EXPECTED[gname]
        g = groups.get(gname, [])  # 已按 left 升序（即"从左往右数"）
        ordinal = ["第1个", "第2个", "第3个"]
        pts = []  # (point_desc, passed)
        # 点①：三个标签均位于该目录项文字下方
        below = len(g) >= 3 and all(to_cm(l.top) >= DIR_TOP_TEXT[gname] for l in g[:3])
        pts.append((f'位于"{gname}"目录项下方', below))
        # 点②③④：从左往右数第N个标签文字 == 指定值
        for i in range(3):
            actual = normalize(g[i].text_frame.text) if i < len(g) else ""
            pts.append((f'从左往右数{ordinal[i]}标签文字为"{exp[i]}"', actual == normalize(exp[i])))
        all_ok = all(p for _, p in pts)
        if all_ok:
            hits.append((+3, RULE_TEXT[gname], "全部点命中"))
        else:
            miss = [d for d, p in pts if not p]
            hits.append((0, RULE_TEXT[gname], "未踩到：" + "；".join(miss)))

    # ---------- 加分项 5：标签整体位置 ----------
    # 细则原文：四组标签分别位于对应目录项文字下方的空白区域，不进入右侧图片区域
    #           或页面底部页码区域。每个目录项下方恰好放置3个标签，同组标签从左至右排列。
    # 逐点踩：
    #   ① 四组标签分别位于"对应目录项文字"下方
    #   ② 不进入右侧图片区域
    #   ③ 不进入页面底部页码区域
    #   ④ 每个目录项下方恰好放置3个标签
    #   ⑤ 同组标签从左至右排列
    # 取各目录项文字形状的真实位置（不硬编码）：通过组名文字匹配 slide2 中的目录项 Text 形状。
    dir_item_shape = {}   # 组名 -> 目录项文字形状
    img_left = None
    page_no_shape = None
    for sh in slide2.shapes:
        if is_rounded_label(sh):
            continue
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img_left = sh.left if img_left is None else min(img_left, sh.left)
        if sh.has_text_frame:
            t = normalize(sh.text_frame.text)
            if t in GROUP_ORDER:
                dir_item_shape[t] = sh
            if t == "02":   # 第2页页码
                page_no_shape = sh

    pos_pts = []  # (point_desc, passed)

    # 点①：四组标签分别位于"对应目录项文字"下方（标签 top >= 该目录项文字 top）
    p1 = True
    for gname in GROUP_ORDER:
        g = groups.get(gname, [])
        ds = dir_item_shape.get(gname)
        if ds is None or not g:
            p1 = False; continue
        if not all((l.top or 0) >= ds.top for l in g):
            p1 = False
    pos_pts.append(("四组标签分别位于对应目录项文字下方", p1))

    # 点②：不进入右侧图片区域（标签右边界 <= 最靠左图片的左边界）
    if img_left is not None:
        p2 = all((l.left or 0) + (l.width or 0) <= img_left for l in labels)
    else:
        p2 = True
    pos_pts.append(("不进入右侧图片区域", p2))

    # 点③：不进入页面底部页码区域（标签下边界 <= 页码 top）
    if page_no_shape is not None:
        p3 = all((l.top or 0) + (l.height or 0) <= page_no_shape.top for l in labels)
    else:
        p3 = True
    pos_pts.append(("不进入页面底部页码区域", p3))

    # 点④：每个目录项下方恰好放置3个标签
    p4 = all(len(groups.get(gname, [])) == 3 for gname in GROUP_ORDER)
    pos_pts.append(("每个目录项下方恰好放置3个标签", p4))

    # 点⑤：同组标签从左至右排列（组内 left 严格递增）
    p5 = True
    for gname in GROUP_ORDER:
        lefts = [l.left or 0 for l in groups.get(gname, [])]
        if lefts != sorted(lefts) or len(set(lefts)) != len(lefts):
            p5 = False
    pos_pts.append(("同组标签从左至右排列", p5))

    pos_ok = all(p for _, p in pos_pts)
    if pos_ok:
        hits.append((+3, RULE_TEXT["位置"], "全部点命中"))
    else:
        miss = [d for d, p in pos_pts if not p]
        hits.append((0, RULE_TEXT["位置"], "未踩到：" + "；".join(miss)))

    # ---------- 加分项 6：尺寸（宽1.8~2.5cm 高0.55~0.85cm；同/不同组大小一致，宽高偏差≤0.1cm；水平间距相等 0.25~0.45cm，间距差≤0.1cm） ----------
    widths = [to_cm(l.width) for l in labels]
    heights = [to_cm(l.height) for l in labels]
    # 细则原文：宽度约1.8至2.5厘米、高度约0.55至0.85厘米，同组及不同组标签大小一致，
    #           宽高偏差不超过0.1厘米。水平间距相等，建议保持0.25至0.45厘米，
    #           任意两个间距差不超过0.1厘米。
    # 逐点踩：
    #   ① 宽度约1.8~2.5厘米（全部标签）
    #   ② 高度约0.55~0.85厘米（全部标签）
    #   ③ 同组及不同组标签大小一致，宽高偏差不超过0.1厘米（即全部标签间宽偏差≤0.1且高偏差≤0.1）
    #   ④ 水平间距相等（任意两个间距差不超过0.1厘米）
    #   ⑤ 建议保持0.25至0.45厘米（"建议"项，仅作提示，不影响是否加分）
    size_pts = []   # (point_desc, passed)

    # 点①：宽度约1.8~2.5厘米
    p1 = all(1.8 <= w <= 2.5 for w in widths)
    size_pts.append(("全部标签宽度约1.8~2.5厘米", p1))

    # 点②：高度约0.55~0.85厘米
    p2 = all(0.55 <= h <= 0.85 for h in heights)
    size_pts.append(("全部标签高度约0.55~0.85厘米", p2))

    # 点③：同组及不同组标签大小一致，宽高偏差不超过0.1厘米
    p3 = bool(widths) and (max(widths) - min(widths) <= 0.1) and (max(heights) - min(heights) <= 0.1)
    size_pts.append(("同组及不同组标签大小一致（宽高偏差≤0.1厘米）", p3))

    # 计算各组内相邻标签的水平间距
    gaps = []
    for gname in GROUP_ORDER:
        g = groups.get(gname, [])
        for a, b in zip(g, g[1:]):
            gaps.append(to_cm(b.left) - (to_cm(a.left) + to_cm(a.width)))

    # 点④：水平间距相等，任意两个间距差不超过0.1厘米
    p4 = bool(gaps) and (max(gaps) - min(gaps) <= 0.1)
    size_pts.append(("水平间距相等（任意两间距差≤0.1厘米）", p4))

    # 是否加分：①②③④四个硬性点全部满足
    size_ok = all(p for _, p in size_pts)

    # 点⑤：建议保持0.25~0.45厘米（"建议"项，仅提示，不参与加分判定）
    p5 = bool(gaps) and all(0.25 <= round(gp, 3) <= 0.45 for gp in gaps)

    if size_ok:
        suggest = "符合" if p5 else "未落入0.25~0.45建议区间(不影响加分)"
        hits.append((+3, RULE_TEXT["尺寸"],
                     f"间距建议项：{suggest} gaps={[round(g,3) for g in gaps]}"))
    else:
        miss = [d for d, p in size_pts if not p]
        hits.append((0, RULE_TEXT["尺寸"], "未踩到：" + "；".join(miss)))

    # ---------- 加分项 7：文字字体 MiSans、字号 10.5磅、水平+垂直居中 ----------
    # 细则原文：第2页全部标签文字：字体统一为MiSans，字号统一为10.5磅，
    #           文字水平居中、垂直居中。
    # 逐点踩（"全部标签"=12个均需满足，"统一"=取值一致）：
    #   ① 字体统一为 MiSans
    #   ② 字号统一为 10.5 磅
    #   ③ 文字水平居中
    #   ④ 文字垂直居中
    #
    # 【字体判定的正确口径】OOXML 中一个 run 的字体分三个槽位：
    #   <a:latin> 西文/数字、<a:ea> 东亚(中文)、<a:cs> 复杂文种。
    #   标签文字是中文，决定中文显示的是 <a:ea>，与 <a:latin> 无关。
    #   因此判定"中文字体"必须读 <a:ea typeface="…"> 子元素：
    #     · 注意 rPr 上的 eastAsia="…" 属性是非法写法，PowerPoint/WPS 不认，不能采信；
    #     · 若 run 含中文却没有合法 <a:ea>，则中文回落到主题/默认东亚字体
    #       （本主题 ea 为空 -> 实际显示为等线类回落字体），绝不能拿 <a:latin> 的值冒充。
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    fonts_seen = set()      # 收集每个 run "实际生效的中文字体"（用于"统一为MiSans"）
    sizes_seen = set()      # 收集每个 run 的字号（用于"统一为10.5磅"）
    halign_ok = True        # 全部段落水平居中
    valign_ok = True        # 全部标签垂直居中

    def _has_cjk(s):
        return any('一' <= ch <= '鿿' for ch in (s or ""))

    def _effective_cjk_font(r):
        """返回该 run 中文字符实际生效的字体名。
        仅认合法的 <a:ea> 子元素；缺失则视为回落（返回 None，代表非 MiSans）。"""
        rPr = r._r.find(f'{{{A}}}rPr')
        if rPr is not None:
            ea = rPr.find(f'{{{A}}}ea')          # 合法的东亚字体子元素
            if ea is not None and ea.get('typeface'):
                return ea.get('typeface')
        # 没有合法 <a:ea> -> 中文回落到主题/默认字体（本文件即等线），按"非 MiSans"处理
        return None

    for l in labels:
        tf = l.text_frame
        # 点④：垂直居中（bodyPr anchor="ctr" -> MIDDLE）
        if "MIDDLE" not in str(tf.vertical_anchor):
            valign_ok = False
        for p in tf.paragraphs:
            # 点③：水平居中（pPr algn="ctr" -> CENTER）
            if "CENTER" not in str(p.alignment):
                halign_ok = False
            for r in p.runs:
                # 点①：字体——标签是中文，取中文实际生效字体（<a:ea>）；
                #       含中文但无合法 <a:ea> 即记为回落字体(None)。
                if _has_cjk(r.text):
                    fonts_seen.add(_effective_cjk_font(r))
                else:
                    fonts_seen.add(r.font.name)   # 纯西文/数字 run 看 latin
                # 点②：字号
                sz = r.font.size
                sizes_seen.add(round(sz.pt, 2) if sz is not None else None)

    font_pts = []  # (point_desc, passed)
    # 点①：字体统一为 MiSans —— 所有 run "实际生效字体"只出现 MiSans 这一个值
    #       （含中文却回落的 run 会把 None 混入集合，使其 != {"MiSans"} 而判失败）
    p1 = (fonts_seen == {"MiSans"})
    font_pts.append(("字体统一为 MiSans", p1))
    # 点②：字号统一为 10.5 磅 —— 所有 run 字号只出现 10.5 这一个值
    p2 = (sizes_seen == {10.5})
    font_pts.append(("字号统一为 10.5 磅", p2))
    # 点③：文字水平居中
    font_pts.append(("文字水平居中", halign_ok))
    # 点④：文字垂直居中
    font_pts.append(("文字垂直居中", valign_ok))

    font_ok = all(p for _, p in font_pts)
    if font_ok:
        hits.append((+3, RULE_TEXT["文字"], "全部点命中"))
    else:
        miss = [d for d, p in font_pts if not p]
        hits.append((0, RULE_TEXT["文字"],
                     "未踩到：" + "；".join(miss) + f"  [fonts={fonts_seen} sizes={sizes_seen}]"))

    # ---------- 加分项 8：标签形状 ----------
    # 细则原文：第2页全部标签形状：使用胶囊形或高圆角矩形，无明显尖角，
    #           边框为无轮廓或与填充色一致。形状使用深绿色填充，文字使用白色，
    #           颜色在12个标签中保持一致。
    # 逐点踩（"全部标签"=12个均需满足）：
    #   ① 使用胶囊形或高圆角矩形
    #   ② 无明显尖角
    #   ③ 边框为无轮廓 或 与填充色一致
    #   ④ 形状使用深绿色填充
    #   ⑤ 文字使用白色
    #   ⑥ 颜色在12个标签中保持一致（填充色一致 且 文字色一致）
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST

    fills = []         # 每个标签的填充色
    txt_colors = []    # 每个 run 的文字色
    p1 = True          # 胶囊形/高圆角矩形
    p2 = True          # 无明显尖角
    p3 = True          # 边框无轮廓或与填充一致

    for l in labels:
        # 点①：胶囊形(PILL/ROUND) 或 高圆角矩形(ROUNDED_RECTANGLE)
        try:
            ast = l.auto_shape_type
        except Exception:
            ast = None
        is_capsule_or_round = ast in (
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGULAR_CALLOUT,
        )
        # 胶囊形在 PPT 中本质是圆角比例很大的圆角矩形；这里把圆角矩形纳入
        if not is_capsule_or_round:
            p1 = False

        # 点②：无明显尖角（圆角矩形的圆角调整值需足够大；普通矩形/无圆角=尖角）
        try:
            adj = list(l.adjustments)
            if not adj or adj[0] < 0.3:   # 圆角比例过小 => 尖角明显
                p2 = False
        except Exception:
            p2 = False

        # 填充色
        try:
            fills.append(str(l.fill.fore_color.rgb))
        except Exception:
            fills.append(None)

        # 文字色
        for p in l.text_frame.paragraphs:
            for r in p.runs:
                try:
                    txt_colors.append(str(r.font.color.rgb))
                except Exception:
                    txt_colors.append(None)

        # 点③：边框为无轮廓 或 与填充色一致
        try:
            ln = l.line
            if ln.fill.type is None:
                pass  # 无轮廓 => 满足
            else:
                lc = None
                try:
                    lc = str(ln.color.rgb)
                except Exception:
                    lc = None
                # 有轮廓时须与填充色一致
                if lc is None or fills[-1] is None or not _close_color(lc, fills[-1]):
                    p3 = False
        except Exception:
            p3 = False

    # 点④：形状使用深绿色填充（全部标签）
    def is_dark_green(hexc):
        if not hexc or len(hexc) != 6:
            return False
        r = int(hexc[0:2], 16); g = int(hexc[2:4], 16); b = int(hexc[4:6], 16)
        return g > r and g > b and g < 170 and r < 120 and b < 140
    p4 = bool(fills) and all(is_dark_green(c) for c in fills)

    # 点⑤：文字使用白色（全部 run）
    p5 = bool(txt_colors) and all(c and c.upper() == "FFFFFF" for c in txt_colors)

    # 点⑥：颜色在12个标签中保持一致（填充色一致 且 文字色一致）
    p6 = (len(set(fills)) == 1) and (len(set(txt_colors)) == 1)

    shape_pts = [
        ("使用胶囊形或高圆角矩形", p1),
        ("无明显尖角", p2),
        ("边框为无轮廓或与填充色一致", p3),
        ("形状使用深绿色填充", p4),
        ("文字使用白色", p5),
        ("颜色在12个标签中保持一致", p6),
    ]
    shape_ok = all(p for _, p in shape_pts)
    if shape_ok:
        hits.append((+3, RULE_TEXT["形状"], "全部点命中"))
    else:
        miss = [d for d, p in shape_pts if not p]
        hits.append((0, RULE_TEXT["形状"],
                     "未踩到：" + "；".join(miss) + f"  [fills={set(fills)} txt={set(txt_colors)}]"))

    # =========================================================
    # 扣分项
    # =========================================================
    # -3：任意三个以上标签大小差异超过1厘米。
    # 细则单一点（扣分项，命中即扣）：存在"三个以上"标签，其大小差异超过1厘米。
    # 判定口径：统计有多少个标签"卷入"了大小差异——
    #   只要某标签与任意另一标签的宽差或高差 > 1cm，即视为该标签卷入差异。
    #   当卷入差异的标签数 ≥ 3（即"三个以上"，含3个）时扣 -3。
    involved = set()
    for i in range(len(widths)):
        for j in range(i + 1, len(widths)):
            if abs(widths[i] - widths[j]) > 1.0 or abs(heights[i] - heights[j]) > 1.0:
                involved.add(i); involved.add(j)
    if len(involved) >= 3:
        hits.append((-3, RULE_TEXT["扣-3-尺寸"], f"卷入差异的标签数={len(involved)}（差异阈值>1cm）"))

    # -3：PPT页数不是8页。
    # 细则单一点（扣分项，命中即扣）：PPT 总页数 ≠ 8。
    if len(prs.slides) != 8:
        hits.append((-3, RULE_TEXT["扣-3-页数"], f"页数={len(prs.slides)}"))

    # -1：文件中出现批注、临时说明文字、红色标记、截图边框或多余占位对象。
    # 细则单一点（扣分项，命中任一类即扣 -1）。范围为"文件中"=整个 PPT 全部页面。
    # 细则点名的 5 类：
    #   ① 批注：PPT 的批注对象（comments / 批注关系），或文字含"批注"
    #   ② 临时说明文字：含临时性说明措辞的文本
    #   ③ 红色标记：红色字体/红色文字标注
    #   ④ 截图边框：以图片形式插入且带边框的截图类对象（图片且 line 有可见轮廓）
    #   ⑤ 多余占位对象：空的/占位符文本（如"占位""placeholder"，或空占位形状）
    junk_detail = []
    JUNK_WORDS = ["批注", "备注", "TODO", "待修改", "说明：", "临时", "placeholder", "占位"]

    # 我们新增的合法标签集合（避免把它们误判为多余对象）
    label_ids = {id(l) for l in labels}

    # ① 批注：检查整个包是否存在批注部件
    try:
        for sld in prs.slides:
            part = sld.part
            for rel in part.rels.values():
                if "comment" in rel.reltype.lower():
                    junk_detail.append("批注对象")
                    break
    except Exception:
        pass

    for sld in prs.slides:
        for sh in sld.shapes:
            if id(sh) in label_ids:
                continue  # 跳过本次新增的合法标签
            # ④ 截图边框：图片且带可见轮廓
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    if sh.line.fill.type is not None and (sh.line.width or 0) > 0:
                        junk_detail.append("截图边框(带轮廓的图片)")
                except Exception:
                    pass
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text
            tn = normalize(t)
            # ②/①(文字)/⑤(占位词)：关键词命中
            for w in JUNK_WORDS:
                if w.lower() in t.lower():
                    junk_detail.append(f"临时/批注/占位文字('{w}')")
            # ⑤ 多余占位对象：占位符类型的形状（is_placeholder）且为空
            try:
                if sh.is_placeholder and tn == "":
                    junk_detail.append("空占位对象")
            except Exception:
                pass
            # ③ 红色标记：红色字体
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    try:
                        c = str(r.font.color.rgb).upper()
                        # 纯红或偏红（R 高、G/B 低）
                        if c == "FF0000" or (len(c) == 6 and
                            int(c[0:2],16) >= 200 and int(c[2:4],16) <= 80 and int(c[4:6],16) <= 80):
                            junk_detail.append("红色标记文字")
                    except Exception:
                        pass

    if junk_detail:
        hits.append((-1, RULE_TEXT["扣-1"],
                     "命中：" + "；".join(sorted(set(junk_detail)))))

    total = sum(s for s, _, _ in hits if s != 0)
    return total, hits


def _close_color(a, b, tol=24):
    try:
        ar=int(a[0:2],16); ag=int(a[2:4],16); ab=int(a[4:6],16)
        br=int(b[0:2],16); bg=int(b[2:4],16); bb=int(b[4:6],16)
        return abs(ar-br)<=tol and abs(ag-bg)<=tol and abs(ab-bb)<=tol
    except Exception:
        return False


# =========================================================
# 主流程：统一入口 evaluate(dir_path)
# =========================================================
SCRIPT_ID = "058"

# 维度2的规则清单：加分项每项满分 +3；扣分项固定负值，命中即扣。
_ADD_RULE_KEYS = ["文化沟通", "工业体系", "全球竞争", "认知偏差",
                  "位置", "尺寸", "文字", "形状"]
_DEDUCT_RULE_SPECS = [
    ("扣-3-尺寸", -3),
    ("扣-3-页数", -3),
    ("扣-1", -1),
]


def _build_dim2_items(hits):
    """把内部 hits 列表转成统一约定的 dim2_items 结构。

    - 加分项：max_delta=+3，delta=0/3，hit 反映是否命中；
    - 扣分项：max_delta 为负值，delta=0（未命中）或 = max_delta（命中），
      未命中项也补齐，保证 100 个脚本能对齐评分项矩阵。
    """
    add_index = {RULE_TEXT[k]: 3 for k in _ADD_RULE_KEYS}
    hit_deduct = {}
    items = []
    for s, desc, detail in hits:
        if desc in add_index:
            items.append({
                "rule": desc,
                "max_delta": add_index[desc],
                "delta": s,
                "hit": s > 0,
                "detail": "",
            })
        else:
            hit_deduct[desc] = (s, detail)
    for key, max_delta in _DEDUCT_RULE_SPECS:
        rule_text = RULE_TEXT[key]
        if rule_text in hit_deduct:
            s, detail = hit_deduct[rule_text]
            items.append({
                "rule": rule_text,
                "max_delta": max_delta,
                "delta": s,
                "hit": True,
                "detail": "",
            })
        else:
            items.append({
                "rule": rule_text,
                "max_delta": max_delta,
                "delta": 0,
                "hit": False,
                "detail": "",
            })
    return items


def _find_pptx(dir_path):
    """在 dir_path 中定位被评估的 .pptx 文件，忽略 Office 临时/隐藏文件。"""
    if not os.path.isdir(dir_path):
        return None
    for name in sorted(os.listdir(dir_path)):
        if name.startswith("~$") or name.startswith("."):
            continue
        if name.lower().endswith(".pptx"):
            return os.path.join(dir_path, name)
    return None


def evaluate(dir_path: str) -> dict:
    """评估 dir_path 目录下的 PPTX 文件，返回结构化结果。

    dir_path 应为脚本所在目录的路径，脚本自行在该目录内定位被评估文档。
    返回结构详见《脚本接口差异与统一建议》§2.2。
    """
    # 满分 = 所有加分项 max_delta 之和（扣分项不参与"满分"计算）
    max_score = sum(3 for _ in _ADD_RULE_KEYS)
    result = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": max_score,
    }
    try:
        pptx_path = _find_pptx(dir_path)
        if pptx_path is None:
            result["status"] = "error"
            result["error"] = f"目录 {dir_path} 中未找到 .pptx 文件"
            return result
        result["file_name"] = os.path.basename(pptx_path)

        prs, err = load(pptx_path)
        if err:
            result["status"] = "error"
            result["error"] = err
            result["dim1_reason"] = err
            return result

        d1_ok, d1_reasons, labels = check_dimension1(prs)
        if not d1_ok:
            result["dim1_pass"] = False
            result["dim1_reason"] = "；".join(r for r in d1_reasons if r.startswith("✗"))
            return result

        result["dim1_pass"] = True
        total, hits = check_dimension2(prs, labels)
        result["dim2_items"] = _build_dim2_items(hits)
        result["total_score"] = max(0, total)
        return result
    except Exception as e:
        result["status"] = "error"
        result["error"] = f"{type(e).__name__}: {e}"
        return result


if __name__ == "__main__":
    # 仅用于本地调试：命令行传入目录，或默认使用脚本所在目录
    _reconfigure = getattr(sys.stdout, "reconfigure", None)
    if _reconfigure is not None:
        try:
            _reconfigure(encoding="utf-8")  # 仅调试路径下调整，避免作为 import 副作用
        except Exception:
            pass
    _target = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(_target), ensure_ascii=False, indent=2))
