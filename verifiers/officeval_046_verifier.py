# -*- coding: utf-8 -*-
"""
自动评估脚本：评估 "两页可编辑复刻图.pptx" 的完成度
评估逻辑：
  维度1（可用与可修改性）：不满足则直接判零分
  维度2（完成度评分）：逐条检查得分点，累计分数
"""
import os
import re
import sys
import json

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 脚本编号（用于返回结构中的 id 字段）
SCRIPT_ID = "046"
# 该脚本评估的目标文档名（作为在目录中定位文档的首选名称；找不到时回退到目录内首个 .pptx）
TARGET_DOC_NAME = "两页可编辑复刻图.pptx"


# ============ 结构化返回相关的辅助函数 ============

_RULE_PREFIX_RE = re.compile(r'^\+(\d+):\s*(.*)$')


def _parse_rule_desc(desc):
    """从 '+N: 描述' 中拆出 (满分, 规则名)。"""
    m = _RULE_PREFIX_RE.match(desc)
    if m:
        return int(m.group(1)), m.group(2).strip()
    return 0, desc


def _dim1_reason_from_results(dim1_results):
    """从 dim1_results 里拼出未通过原因（多条用中文分号连接）。"""
    fails = [d for d, ok in dim1_results if not ok]
    return "；".join(fails)


def _build_dim1_fail_result(dim1_results):
    """维度一未通过时的标准返回结构。"""
    return {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": _dim1_reason_from_results(dim1_results),
        "dim2_items": [],
        "total_score": 0,
        "max_score": 0,
    }


def _build_success_result(dim1_results, dim2_results, total_score):
    """维度一通过、维度二逐项汇总后的标准返回结构。"""
    items = []
    max_total = 0
    for score, desc, hit in dim2_results:
        max_delta, rule = _parse_rule_desc(desc)
        items.append({
            "rule": rule,
            "max_delta": max_delta,
            "delta": int(score),
            "hit": bool(hit),
            "detail": "",
        })
        max_total += max_delta
    dim1_pass_all = all(ok for _, ok in dim1_results)
    return {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": dim1_pass_all,
        "dim1_reason": "" if dim1_pass_all else _dim1_reason_from_results(dim1_results),
        "dim2_items": items,
        "total_score": int(total_score),
        "max_score": max_total,
    }


def _locate_document(dir_path):
    """在给定目录里定位被评估的文档，返回 (file_path, file_name) 或 (None, '')。"""
    if not dir_path or not os.path.isdir(dir_path):
        return None, ""
    preferred = os.path.join(dir_path, TARGET_DOC_NAME)
    if os.path.isfile(preferred):
        return preferred, TARGET_DOC_NAME
    try:
        for name in sorted(os.listdir(dir_path)):
            if name.startswith('~$'):
                continue
            if name.lower().endswith('.pptx'):
                return os.path.join(dir_path, name), name
    except OSError:
        return None, ""
    return None, ""

# ============ 工具函数 ============

def cm(emu):
    """EMU转厘米"""
    if emu is None:
        return 0
    return emu / 360000

def pt_from_emu(emu):
    """EMU转磅"""
    if emu is None:
        return 0
    return emu / 12700

def get_shape_prst(shape):
    """获取形状的preset geometry类型"""
    sp = shape._element
    geom = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    if geom is not None:
        return geom.get('prst')
    return None


def _is_slanted_bottom_right_angle_trapezoid(shape):
    """判定 shape 是否为"向下的直角梯形",且最下侧那条边为"从左上到右下"的斜线。

    "直角梯形"这里定义为顶边水平(两个上顶点 Y 接近), 底边为斜线且从左端偏上
    (Y较小)走到右端偏下(Y较大)。标准 preset ``trapezoid`` 是上下对称的等腰
    梯形, 底边永远水平, 无法表达这种斜底边, 因此必须是 ``custGeom``
    (自由绘制多边形), 从 pathLst/path 中读取 moveTo/lnTo 顶点后判定。

    考虑形状 xfrm 的 flipH/flipV; 旋转超过 45° 时(方向不再"向下")直接判否。
    """
    DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    sp = shape._element
    cust = sp.find(f'.//{DML}custGeom')
    if cust is None:
        return False
    path = cust.find(f'{DML}pathLst/{DML}path')
    if path is None:
        return False
    try:
        path_w = int(path.get('w') or 0)
        path_h = int(path.get('h') or 0)
    except (TypeError, ValueError):
        return False
    if path_w <= 0 or path_h <= 0:
        return False

    # 收集 moveTo / lnTo 顶点(忽略贝塞尔控制点; 直角梯形是折线多边形已足够)
    pts = []
    for el in path.iter():
        tag = el.tag.split('}')[-1]
        if tag in ('moveTo', 'lnTo'):
            pt_el = el.find(f'{DML}pt')
            if pt_el is None:
                continue
            try:
                pts.append((int(pt_el.get('x')), int(pt_el.get('y'))))
            except (TypeError, ValueError):
                continue
    if len(pts) < 4:
        return False

    xfrm = sp.find(f'.//{DML}xfrm')
    rot = 0
    flip_h = flip_v = False
    if xfrm is not None:
        try:
            rot = int(xfrm.get('rot') or 0) / 60000
        except (TypeError, ValueError):
            rot = 0
        flip_h = xfrm.get('flipH') == '1'
        flip_v = xfrm.get('flipV') == '1'
    if abs(rot) > 45:
        return False

    def transformed(pt):
        x, y = pt
        if flip_h:
            x = path_w - x
        if flip_v:
            y = path_h - y
        return (x, y)

    pts = [transformed(p) for p in pts]
    ys_sorted = sorted(pts, key=lambda p: p[1])
    top_two = ys_sorted[:2]
    bottom_two = ys_sorted[-2:]

    # 顶边两点 Y 应接近(水平顶边 → 顶部两处为直角)
    if abs(top_two[0][1] - top_two[1][1]) > path_h * 0.05:
        return False

    # 底边端点按 X 排序: 左端点 Y 应显著小于右端点 Y (从左上→右下)
    left_pt, right_pt = sorted(bottom_two, key=lambda p: p[0])
    if right_pt[0] - left_pt[0] <= 0:
        return False
    if (right_pt[1] - left_pt[1]) < path_h * 0.05:
        return False
    return True

# ---- 主题色标准映射 & 通用解析 ----
# 说明：为了让"非主题色的颜色设定"也能通过判定，这里把 schemeClr 解析成 RGB
# 后统一交给下方的 color_is_* 类别函数（灰/蓝/橙/紫…）判断。任何 schemeClr
# 只要经过标准映射 + lumMod/lumOff/tint/shade 变换后落入目标颜色区间即通过，
# 不再要求必须命中某个特定主题槽位或某个具体十六进制。
STANDARD_SCHEME_COLORS = {
    'dk1': '000000', 'tx1': '000000',
    'lt1': 'FFFFFF', 'bg1': 'FFFFFF', 'background1': 'FFFFFF',
    'dk2': '1F3864', 'tx2': '1F3864',
    'lt2': 'E7E6E6', 'bg2': 'E7E6E6', 'background2': 'E7E6E6',
    'accent1': '4472C4',
    'accent2': 'ED7D31',
    'accent3': 'A5A5A5',
    'accent4': 'FFC000',
    'accent5': '5B9BD5',
    'accent6': '70AD47',
    'hlink': '0563C1',
    'folhlink': '954F72',
    'phclr': '808080',
}

def _apply_color_transforms(hex_color, color_elem):
    """在给定基色上应用 DrawingML 的亮度/色相变换（lumMod / lumOff / tint / shade / satMod）。"""
    if not hex_color:
        return None
    try:
        r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    except Exception:
        return hex_color
    for child in color_elem:
        tag = child.tag.split('}')[-1]
        raw = child.get('val')
        if raw is None:
            continue
        try:
            val = int(raw) / 100000.0
        except Exception:
            continue
        if tag == 'lumMod':
            r, g, b = r * val, g * val, b * val
        elif tag == 'lumOff':
            r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
        elif tag == 'tint':
            r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
        elif tag == 'shade':
            r, g, b = r * val, g * val, b * val
    r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
    return f'{r:02X}{g:02X}{b:02X}'

def _resolve_color_element(container):
    """从含 srgbClr / schemeClr 的容器里解析出十六进制 RGB。"""
    if container is None:
        return None
    DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    srgb = container.find(f'.//{DML}srgbClr')
    if srgb is not None:
        return _apply_color_transforms(srgb.get('val'), srgb)
    scheme = container.find(f'.//{DML}schemeClr')
    if scheme is not None:
        base = STANDARD_SCHEME_COLORS.get((scheme.get('val') or '').lower())
        if base:
            return _apply_color_transforms(base, scheme)
    return None

def get_fill_color(shape):
    """获取形状填充颜色(hex字符串)。兼容 srgbClr 与 schemeClr，避免因未映射主题色导致颜色为空。"""
    sp = shape._element
    DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    for fill in sp.findall(f'.//{DML}solidFill'):
        # 跳过属于 ln (线条) 的 solidFill，避免把线条色当填充
        parent = fill.getparent() if hasattr(fill, 'getparent') else None
        if parent is not None and parent.tag.endswith('}ln'):
            continue
        c = _resolve_color_element(fill)
        if c:
            return c
    return None

def get_line_color(shape):
    """获取线条颜色。兼容 srgbClr 与 schemeClr。"""
    sp = shape._element
    DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    ln = sp.find(f'.//{DML}ln')
    if ln is not None:
        c = _resolve_color_element(ln)
        if c:
            return c
    return None

def get_line_width_pt(shape):
    """获取线宽(磅)。未显示设置线宽时按 PPT 默认线宽 1.0 磅处理。"""
    sp = shape._element
    ln = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    if ln is not None:
        w = ln.get('w')
        if w:
            try:
                return int(w) / 12700
            except Exception:
                pass
    return 1.0

def get_line_dash(shape):
    """获取虚线类型"""
    sp = shape._element
    ln = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
    if ln is not None:
        dash = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstDash')
        if dash is not None:
            return dash.get('val')
    return 'solid'

# ---- 箭头判定 ----
# "一体箭头"包括：所有块状箭头预设（rightArrow/downArrow/…）以及
# 带有箭头端点(headEnd/tailEnd)的 line 与直线连接符(straightConnector*/bentConnector*)。
# 用统一的 helper 后，各处判定只关心"是不是箭头 + 方向/位置/颜色"，
# 不再纠结实现形式究竟是 line 还是 connector。
BLOCK_ARROW_PRSTS = {
    'rightArrow', 'leftArrow', 'upArrow', 'downArrow',
    'leftRightArrow', 'upDownArrow', 'quadArrow', 'leftRightUpArrow',
    'bentArrow', 'uturnArrow', 'bentUpArrow',
    'curvedRightArrow', 'curvedLeftArrow', 'curvedUpArrow', 'curvedDownArrow',
    'stripedRightArrow', 'notchedRightArrow',
    'chevron', 'pentagon', 'homePlate',
    'circularArrow', 'leftCircularArrow', 'leftRightCircularArrow',
}

def _line_arrow_ends(shape):
    """返回 line/连接符两端箭头 type 元组 (headEnd_type, tailEnd_type)。"""
    DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
    sp = shape._element
    head_el = sp.find(f'.//{DML}headEnd')
    tail_el = sp.find(f'.//{DML}tailEnd')
    head = (head_el.get('type', 'none') if head_el is not None else 'none') or 'none'
    tail = (tail_el.get('type', 'none') if tail_el is not None else 'none') or 'none'
    return head, tail

def is_line_or_connector(shape):
    """判断是否为 line 或直线/折线连接符（DrawingML 中都是 cxnSp/prstGeom=line/*Connector*）。"""
    prst = get_shape_prst(shape) or ''
    return prst == 'line' or 'onnector' in prst

def has_line_arrowhead(shape):
    """line/连接符至少一端有非 none 的箭头端点。"""
    if not is_line_or_connector(shape):
        return False
    head, tail = _line_arrow_ends(shape)
    return head not in ('none', '') or tail not in ('none', '')

def has_double_arrowhead(shape):
    """line/连接符两端都有非 none 的箭头端点。"""
    if not is_line_or_connector(shape):
        return False
    head, tail = _line_arrow_ends(shape)
    return head not in ('none', '') and tail not in ('none', '')

def is_arrow_shape(shape):
    """是否为"一体箭头"——块状箭头预设，或带箭头端点的 line/直线连接符。"""
    prst = get_shape_prst(shape)
    if prst in BLOCK_ARROW_PRSTS:
        return True
    return has_line_arrowhead(shape)

def get_text(shape):
    """获取形状文本"""
    if hasattr(shape, 'text') and shape.text:
        return shape.text.replace('\n', '').replace('\r', '').strip()
    return ''

def in_range(val, lo, hi):
    """判断值是否在范围内"""
    return lo <= val <= hi

def color_is_dark_blue(c):
    """判断颜色是否为深蓝/灰蓝"""
    if not c:
        return False
    c = c.upper()
    # 常见深蓝色
    dark_blues = ['09376A', '003366', '1A3D6D', '2E4A7A', '003D7A', '0D47A1',
                  '1E3A5F', '2B4C7E', '1F4E79', '002060', '003399', '1B3F8B',
                  '4178C8', '2E69B3', '345B8C', '2F5496']
    if c in dark_blues:
        return True
    try:
        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        return b > r and b > g and (r + g + b) < 400
    except:
        return False

def color_is_blue(c):
    """判断是否为蓝色系"""
    if not c:
        return False
    try:
        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        return b >= r and b >= g and b > 80
    except:
        return False

def color_is_orange(c):
    """判断是否为橙色"""
    if not c:
        return False
    try:
        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        return r > 180 and g > 80 and g < 200 and b < 100
    except:
        return False

def color_is_red(c):
    """判断是否为红色/深红"""
    if not c:
        return False
    try:
        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        return r > 150 and g < 100 and b < 100
    except:
        return False

def color_is_green(c):
    """判断是否为绿色"""
    if not c:
        return False
    try:
        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        return g > r and g > b and g > 80
    except:
        return False

def color_is_purple(c):
    """判断是否为紫色"""
    if not c:
        return False
    try:
        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        return r > 60 and b > 60 and b >= g and r >= g * 0.8
    except:
        return False

def color_is_yellow(c):
    """判断是否为黄色"""
    if not c:
        return False
    try:
        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        return r > 180 and g > 150 and b < 120 and abs(r - g) < 120
    except:
        return False

def color_is_white_or_light(c):
    """判断是否为白色或浅色"""
    if not c:
        return False
    try:
        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        return (r + g + b) > 600
    except:
        return False

def color_is_light_blue(c):
    """判断是否为浅蓝"""
    if not c:
        return False
    try:
        r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        return b >= r and b >= g and (r + g + b) > 500
    except:
        return False

def find_shapes_with_text(shapes, text, partial=True):
    """查找包含指定文本的形状"""
    results = []
    for s in shapes:
        t = get_text(s)
        if partial and text in t:
            results.append(s)
        elif not partial and t == text:
            results.append(s)
    return results

def find_shapes_in_region(shapes, left_lo, left_hi, top_lo, top_hi):
    """查找位于指定区域内的形状"""
    results = []
    for s in shapes:
        l = cm(s.left)
        t = cm(s.top)
        if left_lo <= l <= left_hi and top_lo <= t <= top_hi:
            results.append(s)
    return results

def shapes_in_area_count(shapes, left_lo, left_hi, top_lo, top_hi):
    """统计区域内形状数量"""
    return len(find_shapes_in_region(shapes, left_lo, left_hi, top_lo, top_hi))


# ============ 主评估逻辑 ============

def evaluate(dir_path: str) -> dict:
    # 在指定目录里定位被评估文档
    filepath, file_name = _locate_document(dir_path)
    if not filepath:
        return {
            "id": SCRIPT_ID,
            "file_name": "",
            "status": "error",
            "error": f"目录中未找到可评估的 .pptx 文件: {dir_path}",
            "dim1_pass": False,
            "dim1_reason": "",
            "dim2_items": [],
            "total_score": 0,
            "max_score": 0,
        }

    try:
        return _evaluate_file(filepath, file_name)
    except Exception as e:
        return {
            "id": SCRIPT_ID,
            "file_name": file_name,
            "status": "error",
            "error": f"{type(e).__name__}: {e}",
            "dim1_pass": False,
            "dim1_reason": "",
            "dim2_items": [],
            "total_score": 0,
            "max_score": 0,
        }


def _evaluate_file(filepath, file_name):
    results = []  # (分数, 描述, 是否通过)

    # ---- 维度1检查 ----
    dim1_pass = True
    dim1_results = []

    # 1. 文件格式
    ext = os.path.splitext(filepath)[1].lower()
    if ext != '.pptx':
        dim1_results.append(("文件格式为.pptx", False))
        dim1_pass = False
    else:
        dim1_results.append(("文件格式为.pptx", True))

    # 2. 尝试打开
    try:
        prs = Presentation(filepath)
    except Exception as e:
        dim1_results.append((f"文件可正常打开 (错误: {e})", False))
        dim1_pass = False
        result = _build_dim1_fail_result(dim1_results)
        result["file_name"] = file_name
        return result
    dim1_results.append(("文件可正常打开", True))

    slides = list(prs.slides)

    # 3. 恰好2页
    if len(slides) != 2:
        dim1_results.append((f"包含2页幻灯片 (实际{len(slides)}页)", False))
        dim1_pass = False
    else:
        dim1_results.append(("包含2页幻灯片", True))

    if not dim1_pass:
        result = _build_dim1_fail_result(dim1_results)
        result["file_name"] = file_name
        return result

    if not dim1_pass:
        result = _build_dim1_fail_result(dim1_results)
        result["file_name"] = file_name
        return result

    # ---- 维度2检查 ----
    slide1_shapes = list(slides[0].shapes)
    slide2_shapes = list(slides[1].shapes)
    dim2_results = []
    total_score = 0

    # ===== 第1页评分 =====

    # +5: 三维立体箱体式边框
    def check_s1_3d_box():
        # 正面矩形 + 左/右侧面 + 顶面共同构成3D箱体，按题干尺寸严格检查
        def box_style_ok(shape):
            DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

            def theme_color_to_rgb(val, is_line):
                val = (val or '').lower()
                if is_line:
                    theme_map = {
                        'dk1': '1F1F1F', 'tx1': '1F1F1F',
                        'dk2': '1F4E79', 'accent1': '4472C4',
                        'accent5': '5B9BD5',
                    }
                else:
                    theme_map = {
                        'lt1': 'FFFFFF', 'bg1': 'FFFFFF', 'background1': 'FFFFFF',
                        'lt2': 'E8EFF7', 'bg2': 'E8EFF7', 'background2': 'E8EFF7',
                    }
                return theme_map.get(val) or STANDARD_SCHEME_COLORS.get(val)

            def office_line_color(shape):
                ln = shape._element.find(f'.//{DML}ln')
                if ln is None:
                    return None
                srgb = ln.find(f'.//{DML}srgbClr')
                if srgb is not None:
                    return srgb.get('val')
                scheme = ln.find(f'.//{DML}schemeClr')
                if scheme is not None:
                    return theme_color_to_rgb(scheme.get('val'), True)
                return None

            def office_fill_color(shape):
                sp = shape._element
                fills = sp.findall(f'.//{DML}solidFill')
                for fill in fills:
                    parent = fill.getparent() if hasattr(fill, 'getparent') else None
                    if parent is not None and parent.tag.endswith('ln'):
                        continue
                    srgb = fill.find(f'.//{DML}srgbClr')
                    if srgb is not None:
                        return srgb.get('val')
                    scheme = fill.find(f'.//{DML}schemeClr')
                    if scheme is not None:
                        mapped = theme_color_to_rgb(scheme.get('val'), False)
                        if mapped:
                            return mapped
                return get_fill_color(shape)

            def color_is_dark_blue_or_gray_blue(c):
                if color_is_dark_blue(c):
                    return True
                if not c:
                    return False
                try:
                    r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                    return (b >= r and b >= g
                            and abs(r - g) <= 35
                            and (b - r) <= 90
                            and 230 <= (r + g + b) <= 560)
                except:
                    return False

            lc = office_line_color(shape)
            lw = get_line_width_pt(shape)
            fc = office_fill_color(shape)
            return (color_is_dark_blue_or_gray_blue(lc)
                    and lw is not None and 0.75 <= lw <= 1.0
                    and (color_is_white_or_light(fc) or color_is_light_blue(fc)))

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        fronts = []
        for s in slide1_shapes:
            l, t, w, h = cm(s.left), cm(s.top), cm(s.width), cm(s.height)
            if (get_shape_prst(s) == 'rect'
                    and 16 <= w <= 18 and 12 <= h <= 14
                    and 3 <= l and shape_right(s) <= 27
                    and 0.8 <= t and shape_bottom(s) <= 18.5
                    and box_style_ok(s)):
                fronts.append(s)

        for front in fronts:
            fl, ft = cm(front.left), cm(front.top)
            fw, fh = cm(front.width), cm(front.height)
            fr, fb = fl + fw, ft + fh

            left_side = None
            right_side = None
            top_face = None

            for s in slide1_shapes:
                if s == front or not box_style_ok(s):
                    continue
                l, t, w, h = cm(s.left), cm(s.top), cm(s.width), cm(s.height)
                r, b = l + w, t + h
                prst = get_shape_prst(s)

                # 左侧面宽度2.2-2.7cm，贴近正面左边，高度与正面基本一致
                if (prst in ('parallelogram', 'freeform', 'rect')
                        and 2.2 <= w <= 2.7
                        and abs(r - fl) <= 0.35
                        and abs(t - ft) <= 0.6 and abs(h - fh) <= 1.0):
                    left_side = s

                # 右侧面宽度3.4-3.9cm，贴近正面右边，高度与正面基本一致
                if (prst in ('parallelogram', 'freeform', 'rect')
                        and 3.4 <= w <= 3.9
                        and abs(l - fr) <= 0.35
                        and abs(t - ft) <= 0.6 and abs(h - fh) <= 1.0):
                    right_side = s

                # 顶面高度1.5-2cm，位于正面上方，横向覆盖正面并延伸到侧面
                if (prst in ('parallelogram', 'freeform', 'rect')
                        and 1.5 <= h <= 2.0
                        and abs(b - ft) <= 0.6
                        and l <= fl + 0.35 and r >= fr - 0.35):
                    top_face = s

            if not (left_side and right_side and top_face):
                continue

            overall_left = min(cm(s.left) for s in (front, left_side, right_side, top_face))
            overall_top = min(cm(s.top) for s in (front, left_side, right_side, top_face))
            overall_right = max(shape_right(s) for s in (front, left_side, right_side, top_face))
            overall_bottom = max(shape_bottom(s) for s in (front, left_side, right_side, top_face))
            if 3 <= overall_left and overall_right <= 27 and 0.8 <= overall_top and overall_bottom <= 18.5:
                return True

        return False

    score, desc = (5, "+5: 第1页三维立体箱体式边框") if check_s1_3d_box() else (0, "+5: 第1页三维立体箱体式边框")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +5: 顶部4个平行四边形分栏标签
    def check_s1_parallelogram_tabs():
        tabs_text = ["运行系统", "膜材料", "污染层生长", "处理系统"]
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5', 'hlink': '0563C1',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF', 'background1': 'FFFFFF',
                'lt2': 'E7E6E6', 'bg2': 'E7E6E6', 'background2': 'E7E6E6',
                'accent3': 'A5A5A5', 'accent6': '70AD47',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_color_transforms(color, color_element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in color_element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_solid_fill(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_color_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_color_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_solid_fill(ln) or get_line_color(shape)

        def office_fill_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_solid_fill(sp_pr) or get_fill_color(shape)

        def color_is_light_gray(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return abs(r - g) <= 25 and abs(g - b) <= 25 and (r + g + b) >= 570
            except:
                return False

        def color_is_light_green(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return g >= r and g >= b and g >= 150 and (r + g + b) >= 500 and not color_is_light_gray(c)
            except:
                return False

        def color_is_light_blue_fill(c):
            return color_is_light_blue(c) and not color_is_light_gray(c) and not color_is_light_green(c)

        expected_fills = [
            color_is_light_blue_fill,
            color_is_light_gray,
            color_is_light_green,
            color_is_light_blue_fill,
        ]

        def parallelogram_style_ok(shape):
            lw = get_line_width_pt(shape)
            return (get_shape_prst(shape) == 'parallelogram'
                    and 4.0 <= cm(shape.width) <= 5.0
                    and 1.5 <= cm(shape.height) <= 2.0
                    and lw is not None and 0.75 <= lw <= 1.25
                    and color_is_dark_blue(office_line_color(shape)))

        def xml_bool(value):
            return value in ('1', 'true', 'True')

        def run_rpr(run):
            return getattr(run._r, 'rPr', None) or run._r.find(f'{DML}rPr')

        def paragraph_def_rpr(para):
            p_pr = getattr(para._p, 'pPr', None) or para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def run_font_names(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (run_rpr(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    font_el = rpr.find(f'{DML}{tag}')
                    if font_el is not None and font_el.get('typeface'):
                        names.add(font_el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(name in ('黑体', 'SimHei', '微软雅黑', 'Microsoft YaHei') for name in names)

        def run_size_pt(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (run_rpr(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def run_bool_style(run, def_rpr, attr, font_attr):
            try:
                if getattr(run.font, font_attr) is True:
                    return True
            except:
                pass
            for rpr in (run_rpr(run), def_rpr):
                if rpr is not None and xml_bool(rpr.get(attr)):
                    return True
            return False

        def run_color(run, def_rpr):
            try:
                if run.font.color.rgb:
                    return str(run.font.color.rgb)
            except:
                pass
            return color_from_solid_fill(run_rpr(run)) or color_from_solid_fill(def_rpr)

        def text_style_ok(shape):
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    def_rpr = paragraph_def_rpr(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        size_pt = run_size_pt(run, def_rpr)
                        if not (font_name_ok(run_font_names(run, def_rpr))
                                and size_pt is not None and 16 <= size_pt <= 18
                                and run_bool_style(run, def_rpr, 'b', 'bold')
                                and run_bool_style(run, def_rpr, 'i', 'italic')
                                and color_is_dark_blue(run_color(run, def_rpr))):
                            return False
                return has_run
            except:
                return False

        def bounds_overlap(a, b, tolerance=0.5):
            al, at, ar, ab = cm(a.left), cm(a.top), shape_right(a), shape_bottom(a)
            bl, bt, br, bb = cm(b.left), cm(b.top), shape_right(b), shape_bottom(b)
            return not (ar < bl - tolerance or br < al - tolerance or ab < bt - tolerance or bb < at - tolerance)

        def text_shape_for_tab(tab, expected_text):
            for shape in slide1_shapes:
                if get_text(shape) == expected_text and bounds_overlap(tab, shape):
                    if text_style_ok(shape):
                        return shape
            return None

        candidates = sorted([s for s in slide1_shapes if parallelogram_style_ok(s)], key=lambda s: cm(s.left))
        if len(candidates) < 4:
            return False

        for start in range(0, len(candidates) - 3):
            paras = candidates[start:start + 4]
            overall_left = min(cm(s.left) for s in paras)
            overall_top = min(cm(s.top) for s in paras)
            overall_right = max(shape_right(s) for s in paras)
            overall_bottom = max(shape_bottom(s) for s in paras)
            if not (5.0 <= overall_left and overall_right <= 27.0 and 0.8 <= overall_top and overall_bottom <= 2.8):
                continue

            matched = True
            for para, expected_txt, fill_ok in zip(paras, tabs_text, expected_fills):
                if not fill_ok(office_fill_color(para)):
                    matched = False
                    break
                if text_shape_for_tab(para, expected_txt) is None:
                    matched = False
                    break
            if matched:
                return True

        return False

    score, desc = (5, "+5: 顶部4个平行四边形分栏标签") if check_s1_parallelogram_tabs() else (0, "+5: 顶部4个平行四边形分栏标签")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +5: 左侧竖向梯形矩形指标栏
    def check_s1_left_indicator():
        texts = ["跨膜压差", "污染迁移速率", "传质阻力"]
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def within_indicator_area(shape):
            return (2.0 <= cm(shape.left)
                    and 2.8 <= cm(shape.top)
                    and shape_right(shape) <= 5.5
                    and shape_bottom(shape) <= 18.0)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF', 'background1': 'FFFFFF',
                'lt2': 'E8EFF7', 'bg2': 'E8EFF7', 'background2': 'E8EFF7',
                'accent3': 'A5A5A5', 'accent6': '70AD47',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_color_transforms(color, color_element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in color_element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_solid_fill(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_color_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_color_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_solid_fill(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_solid_fill(ln) or get_line_color(shape)

        def color_is_light_blue_white(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return b >= r and b >= g and (b - min(r, g)) >= 5 and (r + g + b) >= 600
            except:
                return False

        def color_is_light_blue_gray(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return (b >= r and b >= g
                        and abs(r - g) <= 45
                        and (b - min(r, g)) <= 90
                        and (r + g + b) >= 560)
            except:
                return False

        def color_is_black_or_dark_blue(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return max(r, g, b) <= 45 or color_is_dark_blue(c)
            except:
                return False

        def run_rpr(run):
            return getattr(run._r, 'rPr', None) or run._r.find(f'{DML}rPr')

        def paragraph_def_rpr(para):
            p_pr = getattr(para._p, 'pPr', None) or para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def run_font_names(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (run_rpr(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    font_el = rpr.find(f'{DML}{tag}')
                    if font_el is not None and font_el.get('typeface'):
                        names.add(font_el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(name in ('Microsoft YaHei', '微软雅黑', 'SimHei', '黑体') for name in names)

        def run_size_pt(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (run_rpr(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def run_color(run, def_rpr):
            try:
                if run.font.color.rgb:
                    return str(run.font.color.rgb)
            except:
                pass
            return color_from_solid_fill(run_rpr(run)) or color_from_solid_fill(def_rpr)

        def text_style_ok(shape):
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    def_rpr = paragraph_def_rpr(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        size_pt = run_size_pt(run, def_rpr)
                        if not (font_name_ok(run_font_names(run, def_rpr))
                                and size_pt is not None and 11 <= size_pt <= 13
                                and color_is_black_or_dark_blue(run_color(run, def_rpr))):
                            return False
                return has_run
            except:
                return False

        def bounds_overlap(a, b, tolerance=0.25):
            al, at, ar, ab = cm(a.left), cm(a.top), shape_right(a), shape_bottom(a)
            bl, bt, br, bb = cm(b.left), cm(b.top), shape_right(b), shape_bottom(b)
            return not (ar < bl - tolerance or br < al - tolerance or ab < bt - tolerance or bb < at - tolerance)

        def text_center_y(shape):
            return cm(shape.top) + cm(shape.height) / 2

        indicator_shapes = [s for s in slide1_shapes
                            if within_indicator_area(s)
                            and 2.3 <= cm(s.width) <= 2.8
                            and color_is_light_blue_white(office_fill_color(s))]
        if not indicator_shapes:
            return False

        overall_left = min(cm(s.left) for s in indicator_shapes)
        overall_top = min(cm(s.top) for s in indicator_shapes)
        overall_right = max(shape_right(s) for s in indicator_shapes)
        overall_bottom = max(shape_bottom(s) for s in indicator_shapes)
        if not (2.0 <= overall_left and overall_right <= 5.5 and 2.8 <= overall_top and overall_bottom <= 18.0):
            return False

        text_shapes = []
        for txt in texts:
            matches = [s for s in find_shapes_with_text(slide1_shapes, txt)
                       if get_text(s) == txt and within_indicator_area(s) and text_style_ok(s)]
            if not matches:
                return False
            text_shapes.append(sorted(matches, key=lambda s: cm(s.top))[0])
        if [get_text(s) for s in sorted(text_shapes, key=lambda s: cm(s.top))] != texts:
            return False

        dividers = []
        sorted_texts = sorted(text_shapes, key=lambda s: cm(s.top))
        gaps = [(text_center_y(sorted_texts[0]), text_center_y(sorted_texts[1])),
                (text_center_y(sorted_texts[1]), text_center_y(sorted_texts[2]))]
        for shape in slide1_shapes:
            if not within_indicator_area(shape):
                continue
            lw = get_line_width_pt(shape)
            if not (get_line_dash(shape) in ('dash', 'lgDash', 'sysDash', 'dashDot')
                    and lw is not None and 0.75 <= lw <= 1.25
                    and color_is_light_blue_gray(office_line_color(shape))
                    and abs(cm(shape.height)) <= 0.15):
                continue
            y = cm(shape.top) + abs(cm(shape.height)) / 2
            if any(upper < y < lower for upper, lower in gaps):
                dividers.append(shape)
        if len(dividers) < 2:
            return False

        bottom_text = sorted_texts[2]
        bottom_cells = [s for s in indicator_shapes if bounds_overlap(s, bottom_text, tolerance=0.4)]
        # 传质阻力所在格子必须是"向下的直角梯形",且最下侧线是"从左上到右下"的斜线。
        # 标准 trapezoid 是等腰梯形(底边水平),无法满足"斜底边",这里只认自由形状
        # (custGeom)且顶点满足直角梯形几何特征。
        if not any(_is_slanted_bottom_right_angle_trapezoid(s) for s in bottom_cells):
            return False

        return True

    score, desc = (5, "+5: 左侧竖向指标栏含跨膜压差/污染迁移速率/传质阻力") if check_s1_left_indicator() else (0, "+5: 左侧竖向指标栏含跨膜压差/污染迁移速率/传质阻力")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +5: 左侧指标栏图标
    def check_s1_left_icons():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def bounds(shapes):
            left = min(cm(s.left) for s in shapes)
            top = min(cm(s.top) for s in shapes)
            right = max(shape_right(s) for s in shapes)
            bottom = max(shape_bottom(s) for s in shapes)
            return left, top, right - left, bottom - top

        def size_ok(width, height):
            return 0.9 <= width <= 1.5 and 0.9 <= height <= 1.5

        def in_left_col(shape, extra_top=0.0):
            return (1.5 <= cm(shape.left) <= 4.5
                    and 2.8 + extra_top <= cm(shape.top) <= 18.0)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF', 'background1': 'FFFFFF',
                'lt2': 'E8EFF7', 'bg2': 'E8EFF7',
                'accent3': 'A5A5A5', 'accent6': '70AD47',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def blue_style(shape):
            fc = office_fill(shape)
            lc = office_line(shape)
            return (color_is_blue(fc) or color_is_light_blue(fc) or
                    color_is_blue(lc) or color_is_dark_blue(lc))

        def text_below_group(text_label, icon_group_shapes):
            gl, gt, gw, gh = bounds(icon_group_shapes)
            gb = gt + gh
            for s in slide1_shapes:
                if get_text(s) == text_label:
                    sl, st = cm(s.left), cm(s.top)
                    if gl - 0.5 <= sl <= gl + gw + 0.5 and gb - 0.3 <= st <= gb + 1.5:
                        return True
            return False

        # ---- 1) 圆形指针式压力表图标 ----
        # 主体：圆/弧形，宽高0.9-1.5cm，在左侧栏内，蓝色
        gauge_found = False
        for arc in slide1_shapes:
            if get_shape_prst(arc) not in ('arc', 'ellipse'):
                continue
            aw, ah = cm(arc.width), cm(arc.height)
            if not (size_ok(aw, ah) and in_left_col(arc) and blue_style(arc)):
                continue
            al, at = cm(arc.left), cm(arc.top)
            # 图标内或附近有指针（短线）
            needles = [s for s in slide1_shapes
                       if get_shape_prst(s) == 'line'
                       and al - 0.15 <= cm(s.left) <= al + aw + 0.15
                       and at - 0.15 <= cm(s.top) <= at + ah + 0.15]
            if not needles:
                continue
            if text_below_group("跨膜压差", [arc]):
                gauge_found = True
                break
        if not gauge_found:
            return False

        # ---- 2) 蓝色水滴 + 四个小圆颗粒（一大三小）组合图标 ----
        # 整体宽高0.9-1.5cm；水滴形(teardrop)或椭圆代表水滴。
        # 严格要求四个圆颗粒:1 个显著较大 + 3 个尺寸接近的小圆(尺寸分组)。
        droplet_found = False
        for drop in slide1_shapes:
            if get_shape_prst(drop) not in ('teardrop', 'ellipse'):
                continue
            if not (in_left_col(drop) and blue_style(drop)):
                continue
            dl, dt, dw, dh = cm(drop.left), cm(drop.top), cm(drop.width), cm(drop.height)
            circles = [s for s in slide1_shapes
                       if get_shape_prst(s) == 'ellipse'
                       and s != drop
                       and in_left_col(s)
                       and 0.05 <= cm(s.width) <= 0.55
                       and 0.05 <= cm(s.height) <= 0.55
                       and abs(cm(s.top) - dt) <= 1.8]
            # rubric 明确"一大三小四个",这里严格按四个匹配,并对尺寸做 1+3 分组
            if len(circles) != 4:
                continue
            sizes = sorted((cm(c.width) + cm(c.height)) / 2 for c in circles)
            small_sizes, big_size = sizes[:3], sizes[3]
            small_avg = sum(small_sizes) / 3
            # 大圆平均直径需明显 > 小圆(至少 1.4x),3 个小圆尺寸方差小(极差 <= 40% 均值)
            if not (small_avg > 0
                    and big_size >= small_avg * 1.4
                    and (max(small_sizes) - min(small_sizes)) <= small_avg * 0.4):
                continue
            group = [drop] + circles
            _, _, gw, gh = bounds(group)
            if not size_ok(gw, gh):
                continue
            if text_below_group("污染迁移速率", group):
                droplet_found = True
                break
        if not droplet_found:
            return False

        # ---- 3) 正方形顶部实线 + 下方均匀布满小圆颗粒图标 ----
        # rubric: 颗粒需"均匀地布满"下方区域 ⇒ 行/列间距方差小,颗粒覆盖矩形下部
        # 大部分横向与纵向范围。
        resistance_found = False
        for box in slide1_shapes:
            if get_shape_prst(box) != 'rect':
                continue
            bl, bt, bw, bh = cm(box.left), cm(box.top), cm(box.width), cm(box.height)
            if not (size_ok(bw, bh) and abs(bw - bh) <= 0.35 and in_left_col(box)):
                continue
            top_lines = [s for s in slide1_shapes
                         if get_shape_prst(s) == 'line'
                         and bl - 0.15 <= cm(s.left) <= bl + bw + 0.15
                         and bt - 0.15 <= cm(s.top) <= bt + 0.45
                         and abs(cm(s.height)) <= 0.15
                         and abs(cm(s.width)) >= 0.6]
            if not top_lines:
                continue
            particles = [s for s in slide1_shapes
                         if get_shape_prst(s) == 'ellipse'
                         and bl - 0.15 <= cm(s.left) <= bl + bw + 0.15
                         and bt + 0.15 <= cm(s.top) <= bt + bh + 0.1
                         and 0.05 <= cm(s.width) <= 0.35
                         and 0.05 <= cm(s.height) <= 0.35]
            if len(particles) < 6:
                continue

            # 均匀分布判定:
            #  1) 按 Y 中心聚类为行,按 X 中心聚类为列(容差=平均颗粒直径),行列各 >=2
            #  2) 行内相邻 X 间距 / 列内相邻 Y 间距的相对方差 (var / mean^2) <= 0.5
            #  3) 覆盖检查:颗粒 X 跨度 >= 矩形宽度 50%; Y 跨度 >= 顶线以下高度 50%
            centers = [(cm(p.left) + cm(p.width) / 2, cm(p.top) + cm(p.height) / 2)
                       for p in particles]
            avg_diam = sum(cm(p.width) for p in particles) / len(particles)
            tol = max(avg_diam, 0.1)

            rows = []
            for cx, cy in sorted(centers, key=lambda c: c[1]):
                if rows and abs(cy - rows[-1][0][1]) <= tol:
                    rows[-1].append((cx, cy))
                else:
                    rows.append([(cx, cy)])
            cols = []
            for cx, cy in sorted(centers, key=lambda c: c[0]):
                if cols and abs(cx - cols[-1][0][0]) <= tol:
                    cols[-1].append((cx, cy))
                else:
                    cols.append([(cx, cy)])
            if len(rows) < 2 or len(cols) < 2:
                continue

            def _rel_var(gaps):
                if len(gaps) < 2:
                    return 0.0
                mean = sum(gaps) / len(gaps)
                if mean <= 0:
                    return float('inf')
                var = sum((g - mean) ** 2 for g in gaps) / len(gaps)
                return var / (mean ** 2)

            row_gaps = []
            for row in rows:
                xs = sorted(c[0] for c in row)
                row_gaps.extend(xs[i + 1] - xs[i] for i in range(len(xs) - 1))
            col_gaps = []
            for col in cols:
                ys = sorted(c[1] for c in col)
                col_gaps.extend(ys[i + 1] - ys[i] for i in range(len(ys) - 1))
            if _rel_var(row_gaps) > 0.5 or _rel_var(col_gaps) > 0.5:
                continue

            top_line_bottom = max(cm(t.top) + abs(cm(t.height)) for t in top_lines)
            bottom_area_top = max(top_line_bottom, bt)
            bottom_area_height = (bt + bh) - bottom_area_top
            if bottom_area_height <= 0:
                continue
            xs = [c[0] for c in centers]
            ys = [c[1] for c in centers]
            if (max(xs) - min(xs)) < bw * 0.5:
                continue
            if (max(ys) - min(ys)) < bottom_area_height * 0.5:
                continue

            group = [box] + particles
            if text_below_group("传质阻力", group):
                resistance_found = True
                break
        if not resistance_found:
            return False

        # ---- 箭头检查 ----
        # "跨膜压差""污染迁移速率"各自文字下方有一个垂直向下箭头，长1-1.3cm，蓝色线条/浅蓝填充
        for text_label in ("跨膜压差", "污染迁移速率"):
            text_shapes_matched = [s for s in slide1_shapes if get_text(s) == text_label]
            if not text_shapes_matched:
                return False
            ts = sorted(text_shapes_matched, key=lambda s: cm(s.top))[0]
            tl, tt, tw, tb = cm(ts.left), cm(ts.top), cm(ts.width), shape_bottom(ts)
            arrows = [s for s in slide1_shapes
                      if (is_arrow_shape(s) or get_shape_prst(s) is None)
                      and 1.0 <= abs(cm(s.height)) <= 1.3
                      and tl - 0.4 <= cm(s.left) <= tl + tw + 0.4
                      and tb - 0.3 <= cm(s.top) <= tb + 1.0
                      and (color_is_blue(office_fill(s)) or color_is_light_blue(office_fill(s))
                           or color_is_blue(office_line(s)) or color_is_dark_blue(office_line(s)))]
            if not arrows:
                return False

        return True

    score, desc = (5, "+5: 左侧指标栏图标(压力表/水滴/颗粒)") if check_s1_left_icons() else (0, "+5: 左侧指标栏图标(压力表/水滴/颗粒)")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +5: 4个横向圆角矩形(膜组件/污染沉积层/表面相互作用/溶液性质)
    def check_s1_four_rounded_rects():
        texts = ["膜组件", "污染沉积层", "表面相互作用", "溶液性质"]
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF', 'background1': 'FFFFFF',
                'lt2': 'E8EFF7', 'bg2': 'E8EFF7',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def color_is_white_fill(c):
            # 白色填充：RGB 均 >= 240
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r >= 240 and g >= 240 and b >= 240
            except:
                return False

        def rpr_of(run):
            rpr = getattr(run._r, 'rPr', None)
            return rpr if rpr is not None else run._r.find(f'{DML}rPr')

        def def_rpr_of(para):
            p_pr = getattr(para._p, 'pPr', None)
            if p_pr is None:
                p_pr = para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def font_names_of(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    el = rpr.find(f'{DML}{tag}')
                    if el is not None and el.get('typeface'):
                        names.add(el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(n in ('Microsoft YaHei', '微软雅黑', 'SimHei', '黑体') for n in names)

        def size_pt_of(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def is_bold(run, def_rpr):
            try:
                if run.font.bold is True:
                    return True
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('b') in ('1', 'true'):
                    return True
            return False

        def text_style_ok(shape):
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    def_rpr = def_rpr_of(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        sz = size_pt_of(run, def_rpr)
                        if not (font_name_ok(font_names_of(run, def_rpr))
                                and sz is not None and 12 <= sz <= 14
                                and is_bold(run, def_rpr)):
                            return False
                return has_run
            except:
                return False

        def in_range_box(shape):
            return (5.5 <= cm(shape.left)
                    and shape_right(shape) <= 24.0
                    and 3.4 <= cm(shape.top)
                    and shape_bottom(shape) <= 5.0)

        def rounded_rect_ok(shape):
            if get_shape_prst(shape) != 'roundRect':
                return False
            w, h = cm(shape.width), cm(shape.height)
            if not (2.0 <= w <= 3.0 and 1.0 <= h <= 1.5):
                return False
            if not in_range_box(shape):
                return False
            lc = office_line_color(shape)
            lw = get_line_width_pt(shape)
            fc = office_fill(shape)
            return ((color_is_blue(lc) or color_is_dark_blue(lc))
                    and get_line_dash(shape) == 'solid'
                    and lw is not None and 0.75 <= lw <= 1.25
                    and (color_is_white_fill(fc) or fc is None))

        def shapes_overlap(a, b, tol=0.5):
            al, at = cm(a.left), cm(a.top)
            ar, ab = shape_right(a), shape_bottom(a)
            bl, bt = cm(b.left), cm(b.top)
            br, bb = shape_right(b), shape_bottom(b)
            return not (ar < bl - tol or br < al - tol or ab < bt - tol or bb < at - tol)

        rects = sorted([s for s in slide1_shapes if rounded_rect_ok(s)], key=lambda s: cm(s.left))

        found_texts = set()
        used_rects = set()
        for txt in texts:
            for ts in slide1_shapes:
                if get_text(ts) != txt:
                    continue
                if not text_style_ok(ts):
                    continue
                for rect in rects:
                    if id(rect) in used_rects:
                        continue
                    if shapes_overlap(rect, ts):
                        found_texts.add(txt)
                        used_rects.add(id(rect))
                        break
                if txt in found_texts:
                    break

        if len(found_texts) < 4:
            return False

        matched_rects = [s for s in rects if id(s) in used_rects]
        overall_left = min(cm(s.left) for s in matched_rects)
        overall_right = max(shape_right(s) for s in matched_rects)
        overall_top = min(cm(s.top) for s in matched_rects)
        overall_bottom = max(shape_bottom(s) for s in matched_rects)
        return (5.5 <= overall_left and overall_right <= 24.0
                and 3.4 <= overall_top and overall_bottom <= 5.0)

    score, desc = (5, "+5: 4个横向圆角矩形(膜组件等)") if check_s1_four_rounded_rects() else (0, "+5: 4个横向圆角矩形(膜组件等)")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: 4个圆角矩形间3个向右箭头
    def check_s1_arrows_between():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_center_y(shape):
            return cm(shape.top) + cm(shape.height) / 2

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def office_line_width_pt(shape):
            # 先尝试从 spPr/ln 读取，回退到公共函数
            sp_pr = shape._element.find(f'.//{DML}spPr')
            if sp_pr is not None:
                ln = sp_pr.find(f'{DML}ln')
                if ln is not None:
                    w = ln.get('w')
                    if w:
                        return int(w) / 12700
            return get_line_width_pt(shape)

        def has_arrowhead(shape):
            sp = shape._element
            for end_tag in ('headEnd', 'tailEnd'):
                end_el = sp.find(f'.//{DML}{end_tag}')
                if end_el is not None:
                    arrow_type = end_el.get('type', 'none')
                    if arrow_type not in ('none', ''):
                        return True
            return False

        def is_rightward_arrow(shape):
            prst = get_shape_prst(shape)
            if prst in ('rightArrow', 'notchedRightArrow', 'pentagon', 'chevron',
                        'stripedRightArrow', 'homePlate', 'bentArrow'):
                return True
            # line 或直线/折线连接符：横向为主 + 至少一端有箭头端点
            if is_line_or_connector(shape):
                w, h = cm(shape.width), abs(cm(shape.height))
                if w < 0:
                    return False
                if h > w * 0.5:
                    return False
                return has_line_arrowhead(shape)
            return False

        # 先找出4个圆角矩形（从 check_s1_four_rounded_rects 同样条件）
        rect_area_shapes = [s for s in slide1_shapes
                            if get_shape_prst(s) == 'roundRect'
                            and 2.0 <= cm(s.width) <= 3.0
                            and 1.0 <= cm(s.height) <= 1.5
                            and 5.5 <= cm(s.left)
                            and shape_right(s) <= 24.0
                            and 3.4 <= cm(s.top)
                            and cm(s.top) + cm(s.height) <= 5.0]
        if len(rect_area_shapes) < 4:
            # 圆角矩形不足时退化到纯区域检查
            rect_center_y_mid = 4.2
            rect_center_y_range = 1.0
        else:
            rect_area_shapes_sorted = sorted(rect_area_shapes, key=lambda s: cm(s.left))[:4]
            all_cy = [shape_center_y(s) for s in rect_area_shapes_sorted]
            rect_center_y_mid = sum(all_cy) / len(all_cy)
            rect_center_y_range = 0.6  # 箭头纵向中心在矩形中心 ±0.6cm 内算对齐

        found_arrows = []
        for s in slide1_shapes:
            if not is_rightward_arrow(s):
                continue
            w = cm(s.width)
            if not (0.8 <= w <= 1.2):
                continue
            lw = office_line_width_pt(s)
            if lw is None or not (1.0 <= lw <= 1.5):
                continue
            lc = office_line_color(s)
            if not (color_is_blue(lc) or color_is_dark_blue(lc)):
                continue
            if get_line_dash(s) not in ('solid', None):
                continue
            if not (5.5 <= cm(s.left) and shape_right(s) <= 24.0):
                continue
            arrow_cy = cm(s.top) + abs(cm(s.height)) / 2
            if abs(arrow_cy - rect_center_y_mid) <= rect_center_y_range:
                found_arrows.append(s)

        return len(found_arrows) >= 3

    score, desc = (3, "+3: 圆角矩形间3个向右箭头") if check_s1_arrows_between() else (0, "+3: 圆角矩形间3个向右箭头")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: 蓝色虚线连接结构(伞状汇入)
    def check_s1_dashed_lines():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def office_line_width_pt(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            if sp_pr is not None:
                ln = sp_pr.find(f'{DML}ln')
                if ln is not None:
                    w = ln.get('w')
                    if w:
                        return int(w) / 12700
            return get_line_width_pt(shape)

        def color_is_blue_or_gray_blue(c):
            if not c:
                return False
            if color_is_blue(c) or color_is_dark_blue(c):
                return True
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                # 灰蓝：蓝色通道最大，三通道差距较小，整体偏中等亮度
                return (b >= r and b >= g
                        and abs(r - g) <= 40
                        and (b - min(r, g)) <= 80
                        and 200 <= (r + g + b) <= 600)
            except:
                return False

        # 虚线类型集合（排除实线和None）
        dash_types = {'dash', 'lgDash', 'sysDash', 'dashDot', 'lgDashDot',
                      'sysDashDot', 'sysDashDotDot', 'dot'}

        # 收集所有符合条件的虚线（整体位于 top 4.6–6.4cm 范围内）
        dashes = []
        for s in slide1_shapes:
            t, b = cm(s.top), shape_bottom(s)
            # 整体纵向范围在 4.6–6.4 内
            if not (4.6 <= t and b <= 6.4):
                continue
            if get_line_dash(s) not in dash_types:
                continue
            lw = office_line_width_pt(s)
            if lw is None or not (0.75 <= lw <= 1.0):
                continue
            lc = office_line_color(s)
            if not color_is_blue_or_gray_blue(lc):
                continue
            dashes.append(s)

        if len(dashes) < 4:
            return False

        # 伞状结构验证：虚线应覆盖横向矩形区域（左侧范围5–24cm），
        # 且存在向右收窄汇聚到中部的结构
        # 检查：有虚线起点或终点分布在横向较宽范围（≥10cm跨度），
        # 且有虚线终点落在中部区域（横向10–18cm）
        lefts = sorted(cm(s.left) for s in dashes)
        rights = sorted(shape_right(s) for s in dashes)
        overall_span = max(rights) - min(lefts)

        # 横向跨度需覆盖足够范围（至少10cm，体现伞形覆盖4个矩形）
        if overall_span < 10.0:
            return False

        # 至少有2条虚线的右端或左端落在中部汇聚区（距幻灯片左8–18cm）
        mid_convergence = sum(
            1 for s in dashes
            if 8.0 <= shape_right(s) <= 19.0 or 8.0 <= cm(s.left) <= 19.0
        )
        return mid_convergence >= 2

    score, desc = (3, "+3: 蓝色虚线连接结构") if check_s1_dashed_lines() else (0, "+3: 蓝色虚线连接结构")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 蓝色圆角矩形说明框"膜污染演化..."
    def check_s1_membrane_box():
        expected_text = "膜污染演化是影响通量衰减的重要因素之一"
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF', 'background1': 'FFFFFF',
                'lt2': 'E8EFF7', 'bg2': 'E8EFF7', 'background2': 'E8EFF7',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def office_line_width_pt(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            if sp_pr is not None:
                ln = sp_pr.find(f'{DML}ln')
                if ln is not None:
                    w = ln.get('w')
                    if w:
                        return int(w) / 12700
            return get_line_width_pt(shape)

        def color_is_light_blue_fill(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return b >= r and b >= g and (r + g + b) >= 500
            except:
                return False

        def rpr_of(run):
            rpr = getattr(run._r, 'rPr', None)
            return rpr if rpr is not None else run._r.find(f'{DML}rPr')

        def def_rpr_of(para):
            p_pr = getattr(para._p, 'pPr', None)
            if p_pr is None:
                p_pr = para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def font_names_of(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    el = rpr.find(f'{DML}{tag}')
                    if el is not None and el.get('typeface'):
                        names.add(el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(n in ('Microsoft YaHei', '微软雅黑', 'SimHei', '黑体') for n in names)

        def size_pt_of(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def run_color(run, def_rpr):
            try:
                if run.font.color.rgb:
                    return str(run.font.color.rgb)
            except:
                pass
            return color_from_elem(rpr_of(run)) or color_from_elem(def_rpr)

        def text_style_ok(shape):
            try:
                if get_text(shape) != expected_text:
                    return False
                has_run = False
                for para in shape.text_frame.paragraphs:
                    def_rpr = def_rpr_of(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        size_pt = size_pt_of(run, def_rpr)
                        if not (font_name_ok(font_names_of(run, def_rpr))
                                and size_pt is not None and 13 <= size_pt <= 15
                                and color_is_dark_blue(run_color(run, def_rpr))):
                            return False
                return has_run
            except:
                return False

        for box in slide1_shapes:
            if get_shape_prst(box) != 'roundRect':
                continue
            l, t, w, h = cm(box.left), cm(box.top), cm(box.width), cm(box.height)
            if not (8.0 <= l and shape_right(box) <= 20.0
                    and 5.8 <= t and shape_bottom(box) <= 7.2
                    and 10.5 <= w <= 11.5 and 1.0 <= h <= 1.3):
                continue
            lc = office_line_color(box)
            lw = office_line_width_pt(box)
            fc = office_fill(box)
            if not ((color_is_blue(lc) or color_is_dark_blue(lc))
                    and get_line_dash(box) == 'solid'
                    and lw is not None and 0.75 <= lw <= 1.25
                    and color_is_light_blue_fill(fc)):
                continue
            for text_shape in find_shapes_with_text(slide1_shapes, expected_text, partial=False):
                tl, tt = cm(text_shape.left), cm(text_shape.top)
                tr, tb = shape_right(text_shape), shape_bottom(text_shape)
                if l <= tl and tr <= l + w and t <= tt and tb <= t + h and text_style_ok(text_shape):
                    return True
        return False

    score, desc = (1, "+1: 膜污染演化说明框") if check_s1_membrane_box() else (0, "+1: 膜污染演化说明框")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 说明框下方箭头指向橙色框
    def check_s1_arrow_to_orange():
        def has_valid_orange_box():
            for box in slide1_shapes:
                l, t, w, h = cm(box.left), cm(box.top), cm(box.width), cm(box.height)
                if not (get_shape_prst(box) == 'roundRect'
                        and 8 <= l <= 21 and 7.5 <= t <= 10.2
                        and 12 <= w <= 13 and 1.8 <= h <= 2.2
                        and color_is_orange(get_line_color(box))
                        and get_line_dash(box) in ('dash', 'lgDash', 'sysDash', 'dashDot', 'lgDashDot')
                        and 0.75 <= (get_line_width_pt(box) or 0) <= 1.25):
                    continue
                fc = get_fill_color(box)
                if not fc:
                    continue
                try:
                    r, g, b = int(fc[0:2], 16), int(fc[2:4], 16), int(fc[4:6], 16)
                    if r > 245 and g > 235 and b > 220:
                        return True
                except:
                    pass
            return False

        # 必须建立在合规蓝色说明框和合规橙色虚线框都存在的基础上
        if not check_s1_membrane_box() or not has_valid_orange_box():
            return False
        # 说明框下方必须有深蓝色实线向下箭头，高度1-1.2cm，指向橙色框中心
        for s in slide1_shapes:
            l, t, w, h = cm(s.left), cm(s.top), cm(s.width), cm(s.height)
            # 允许 downArrow / bentArrow 等块状箭头，或带箭头端点的 line / 直线连接符
            if not ((is_arrow_shape(s))
                    and 12.5 <= l <= 14.0 and 7.0 <= t <= 8.5
                    and abs(w) <= 0.2 and 1.0 <= abs(h) <= 1.2):
                continue
            if not (get_line_dash(s) == 'solid' and color_is_dark_blue(get_line_color(s))):
                continue
            return True
        return False

    score, desc = (1, "+1: 说明框下方箭头指向橙色框") if check_s1_arrow_to_orange() else (0, "+1: 说明框下方箭头指向橙色框")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 橙色虚线边框圆角矩形
    def check_s1_orange_dashed_box():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            # Office 主题橙色常见映射
            scheme_map = {
                'accent2': 'ED7D31', 'accent6': 'FFC000',
                'dk1': '000000', 'tx1': '000000',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def office_line_width_pt(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            if sp_pr is not None:
                ln = sp_pr.find(f'{DML}ln')
                if ln is not None:
                    w = ln.get('w')
                    if w:
                        return int(w) / 12700
            return get_line_width_pt(shape)

        def color_is_orange_line(c):
            # 橙色：R 最大，G 居中，B 最小；覆盖从纯橙到金黄
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r > 180 and g >= 60 and b < 160 and r > g and r > b
            except:
                return False

        def color_is_very_light_orange(c):
            # 极浅橙色：R 最高，整体接近白色但 R 明显大于 B
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r >= 240 and g >= 220 and b >= 190 and r >= g and r > b
            except:
                return False

        dash_types = {'dash', 'lgDash', 'sysDash', 'dashDot', 'lgDashDot',
                      'sysDashDot', 'sysDashDotDot', 'dot'}

        for s in slide1_shapes:
            if get_shape_prst(s) != 'roundRect':
                continue
            l, t = cm(s.left), cm(s.top)
            w, h = cm(s.width), cm(s.height)
            # 位置：左 8–21cm，上 7.5–10.2cm（整体）
            if not (8.0 <= l and shape_right(s) <= 21.0
                    and 7.5 <= t and shape_bottom(s) <= 10.2):
                continue
            # 尺寸：宽 12–13cm，高 1.8–2.2cm
            if not (12.0 <= w <= 13.0 and 1.8 <= h <= 2.2):
                continue
            lc = office_line_color(s)
            lw = office_line_width_pt(s)
            # 边线橙色虚线，线宽 0.75–1.25 磅
            if not (color_is_orange_line(lc)
                    and get_line_dash(s) in dash_types
                    and lw is not None and 0.75 <= lw <= 1.25):
                continue
            # 填充极浅橙色
            fc = office_fill(s)
            if color_is_very_light_orange(fc):
                return True
        return False

    score, desc = (1, "+1: 橙色虚线边框圆角矩形") if check_s1_orange_dashed_box() else (0, "+1: 橙色虚线边框圆角矩形")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 橙色框内文本
    def check_s1_orange_text():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'accent2': 'ED7D31', 'accent6': 'FFC000',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def rpr_of(run):
            rpr = getattr(run._r, 'rPr', None)
            return rpr if rpr is not None else run._r.find(f'{DML}rPr')

        def def_rpr_of(para):
            p_pr = getattr(para._p, 'pPr', None)
            if p_pr is None:
                p_pr = para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def font_names_of(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    el = rpr.find(f'{DML}{tag}')
                    if el is not None and el.get('typeface'):
                        names.add(el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(n in ('Microsoft YaHei', '微软雅黑', 'SimHei', '黑体') for n in names)

        def size_pt_of(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def is_bold(run, def_rpr):
            try:
                if run.font.bold is True:
                    return True
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('b') in ('1', 'true'):
                    return True
            return False

        def run_color(run, def_rpr):
            try:
                if run.font.color.rgb:
                    return str(run.font.color.rgb)
            except:
                pass
            return color_from_elem(rpr_of(run)) or color_from_elem(def_rpr)

        def color_is_black_or_dark_blue(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                is_black = max(r, g, b) <= 45
                return is_black or color_is_dark_blue(c)
            except:
                return False

        def color_is_orange_or_red_orange(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                # 橙色：R 最大，G 居中，B 较小
                return r > 180 and g >= 50 and b < 160 and r > g and r > b
            except:
                return False

        def para_is_centered(para):
            try:
                from pptx.enum.text import PP_ALIGN
                if para.alignment in (PP_ALIGN.CENTER,):
                    return True
                p_pr = para._p.find(f'{DML}pPr')
                if p_pr is not None and p_pr.get('algn') == 'ctr':
                    return True
            except:
                pass
            return False

        def text_style_ok(shape, expected_text, size_lo, size_hi, color_check, need_bold):
            if get_text(shape) != expected_text:
                return False
            try:
                has_run = False
                centered = False
                for para in shape.text_frame.paragraphs:
                    def_rpr = def_rpr_of(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        sz = size_pt_of(run, def_rpr)
                        if not (font_name_ok(font_names_of(run, def_rpr))
                                and sz is not None and size_lo <= sz <= size_hi
                                and color_check(run_color(run, def_rpr))):
                            return False
                        if need_bold and not is_bold(run, def_rpr):
                            return False
                    if para_is_centered(para):
                        centered = True
                return has_run and centered
            except:
                return False

        # 橙色框范围约束（位置在橙色框内）
        def in_orange_box(shape):
            return (8.0 <= cm(shape.left) and shape_right(shape) <= 21.0
                    and 7.5 <= cm(shape.top) and shape_bottom(shape) <= 10.2)

        t1_shapes = [s for s in slide1_shapes
                     if in_orange_box(s)
                     and text_style_ok(s, "精细调控膜表面结构与运行参数",
                                       14, 16, color_is_black_or_dark_blue, True)]
        t2_shapes = [s for s in slide1_shapes
                     if in_orange_box(s)
                     and text_style_ok(s, "优化污染控制状态",
                                       18, 20, color_is_orange_or_red_orange, True)]

        if not t1_shapes or not t2_shapes:
            return False

        # "精细调控..."在上，"优化污染..."在下
        t1 = sorted(t1_shapes, key=lambda s: cm(s.top))[0]
        t2 = sorted(t2_shapes, key=lambda s: cm(s.top))[0]
        return cm(t1.top) < cm(t2.top)

    score, desc = (1, "+1: 橙色框内文本") if check_s1_orange_text() else (0, "+1: 橙色框内文本")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 橙色框下方连接线(竖线+横线+箭头)
    def check_s1_connector_lines():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def shape_center_x(shape):
            return cm(shape.left) + cm(shape.width) / 2

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def is_blue_color(c):
            return color_is_blue(c) or color_is_dark_blue(c)

        # ---- 1) 橙色框中心下方蓝色竖向实线，长 0.5–0.8cm ----
        vertical_lines = []
        for s in slide1_shapes:
            if get_shape_prst(s) != 'line':
                continue
            w, h = cm(s.width), abs(cm(s.height))
            if not (abs(w) <= 0.2 and 0.5 <= h <= 0.8):
                continue
            if get_line_dash(s) not in ('solid', None):
                continue
            lc = office_line_color(s)
            if not is_blue_color(lc):
                continue
            vertical_lines.append(s)

        if not vertical_lines:
            return False

        # ---- 2) 向下连接蓝色横向实线，长 8–10cm ----
        horizontal_lines = []
        for s in slide1_shapes:
            if get_shape_prst(s) != 'line':
                continue
            w, h = cm(s.width), abs(cm(s.height))
            if not (8.0 <= w <= 10.0 and h <= 0.2):
                continue
            if get_line_dash(s) not in ('solid', None):
                continue
            lc = office_line_color(s)
            if not is_blue_color(lc):
                continue
            horizontal_lines.append(s)

        if not horizontal_lines:
            return False

        # ---- 3) 横线左右两端各有一个竖向实心蓝色箭头，高 0.5–0.7cm ----
        down_arrows = []
        for s in slide1_shapes:
            prst = get_shape_prst(s)
            # downArrow 预设，或带向下箭头端点的竖线/直线连接符
            if prst == 'downArrow':
                h = cm(s.height)
                if not (0.5 <= h <= 0.7):
                    continue
                fc = office_fill(s)
                lc = office_line_color(s)
                if is_blue_color(fc) or is_blue_color(lc):
                    down_arrows.append(s)
            elif is_line_or_connector(s):
                h = abs(cm(s.height))
                w = abs(cm(s.width))
                if not (0.5 <= h <= 0.7 and w <= 0.2):
                    continue
                if not has_line_arrowhead(s):
                    continue
                lc = office_line_color(s)
                if is_blue_color(lc):
                    down_arrows.append(s)

        if len(down_arrows) < 2:
            return False

        # ---- 结构连贯性验证 ----
        # 竖线底端与横线应纵向相近（±0.5cm），且横向在横线覆盖范围内
        vline = sorted(vertical_lines, key=lambda s: cm(s.top))[0]
        hline = sorted(horizontal_lines, key=lambda s: cm(s.left))[0]
        vline_bottom = cm(vline.top) + abs(cm(vline.height))
        hline_top = cm(hline.top)
        if abs(vline_bottom - hline_top) > 0.5:
            return False

        # 两个箭头应分布在横线左右两端附近（左端和右端各一个）
        hline_left = cm(hline.left)
        hline_right = shape_right(hline)
        hline_mid = (hline_left + hline_right) / 2
        arrows_left = [a for a in down_arrows if shape_center_x(a) <= hline_mid]
        arrows_right = [a for a in down_arrows if shape_center_x(a) > hline_mid]
        return bool(arrows_left and arrows_right)

    score, desc = (1, "+1: 橙色框下方连接线结构") if check_s1_connector_lines() else (0, "+1: 橙色框下方连接线结构")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: "界面相容优化"模块
    def check_s1_interface_module():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
                'lt2': 'E8EFF7', 'bg2': 'E8EFF7',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def office_line_width_pt(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            if sp_pr is not None:
                ln = sp_pr.find(f'{DML}ln')
                if ln is not None:
                    w = ln.get('w')
                    if w:
                        return int(w) / 12700
            return get_line_width_pt(shape)

        def color_is_pale_blue(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return b >= r and b >= g and b >= 170 and (r + g + b) >= 520
            except:
                return False

        def rpr_of(run):
            rpr = getattr(run._r, 'rPr', None)
            return rpr if rpr is not None else run._r.find(f'{DML}rPr')

        def def_rpr_of(para):
            p_pr = getattr(para._p, 'pPr', None)
            if p_pr is None:
                p_pr = para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def font_names_of(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    el = rpr.find(f'{DML}{tag}')
                    if el is not None and el.get('typeface'):
                        names.add(el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(n in ('Microsoft YaHei', '微软雅黑', 'SimHei', '黑体') for n in names)

        def size_pt_of(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def is_bold(run, def_rpr):
            try:
                if run.font.bold is True:
                    return True
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('b') in ('1', 'true'):
                    return True
            return False

        def run_color(run, def_rpr):
            try:
                if run.font.color.rgb:
                    return str(run.font.color.rgb)
            except:
                pass
            return color_from_elem(rpr_of(run)) or color_from_elem(def_rpr)

        def title_style_ok(shape):
            if get_text(shape) != "界面相容优化":
                return False
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    def_rpr = def_rpr_of(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        size_pt = size_pt_of(run, def_rpr)
                        if not (font_name_ok(font_names_of(run, def_rpr))
                                and size_pt is not None and 14 <= size_pt <= 16
                                and is_bold(run, def_rpr)
                                and (color_is_blue(run_color(run, def_rpr)) or color_is_dark_blue(run_color(run, def_rpr)))):
                            return False
                return has_run
            except:
                return False

        for box in slide1_shapes:
            if get_shape_prst(box) != 'roundRect':
                continue
            l, t, w, h = cm(box.left), cm(box.top), cm(box.width), cm(box.height)
            if not (5.5 <= l and shape_right(box) <= 14.5
                    and 10.0 <= t and shape_bottom(box) <= 16.0
                    and 7.0 <= w <= 8.0 and 4.0 <= h <= 5.2):
                continue
            lc = office_line_color(box)
            lw = office_line_width_pt(box)
            fc = office_fill(box)
            if not ((color_is_blue(lc) or color_is_dark_blue(lc))
                    and get_line_dash(box) == 'solid'
                    and lw is not None and 0.75 <= lw <= 1.25
                    and color_is_pale_blue(fc)):
                continue
            for title in find_shapes_with_text(slide1_shapes, "界面相容优化", partial=False):
                tl, tt = cm(title.left), cm(title.top)
                tr, tb = shape_right(title), shape_bottom(title)
                if l <= tl and tr <= l + w and t <= tt and tb <= t + h and title_style_ok(title):
                    return True
        return False

    score, desc = (1, "+1: 界面相容优化模块") if check_s1_interface_module() else (0, "+1: 界面相容优化模块")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: 界面相容优化下4个小圆角矩形
    def check_s1_interface_subitems():
        texts = ["增强表面亲水性", "改进孔结构", "膜基材", "表面涂层"]
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
                'lt2': 'E8EFF7', 'bg2': 'E8EFF7',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def office_line_width_pt(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            if sp_pr is not None:
                ln = sp_pr.find(f'{DML}ln')
                if ln is not None:
                    w = ln.get('w')
                    if w:
                        return int(w) / 12700
            return get_line_width_pt(shape)

        def color_is_light_blue_or_white(c):
            if not c:
                return True  # 无填充色视为白色背景
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                is_white = r >= 235 and g >= 235 and b >= 235
                is_light_blue = b >= r and b >= g and b >= 170 and (r + g + b) >= 500
                return is_white or is_light_blue
            except:
                return False

        def color_is_blue_line(c):
            return color_is_blue(c) or color_is_light_blue(c) or color_is_dark_blue(c)

        def rpr_of(run):
            rpr = getattr(run._r, 'rPr', None)
            return rpr if rpr is not None else run._r.find(f'{DML}rPr')

        def def_rpr_of(para):
            p_pr = getattr(para._p, 'pPr', None)
            if p_pr is None:
                p_pr = para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def font_names_of(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    el = rpr.find(f'{DML}{tag}')
                    if el is not None and el.get('typeface'):
                        names.add(el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(n in ('Microsoft YaHei', '微软雅黑', 'SimHei', '黑体') for n in names)

        def size_pt_of(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def text_style_ok(shape):
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    def_rpr = def_rpr_of(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        sz = size_pt_of(run, def_rpr)
                        if not (font_name_ok(font_names_of(run, def_rpr))
                                and sz is not None and 11 <= sz <= 13):
                            return False
                return has_run
            except:
                return False

        def small_roundrect_ok(shape):
            if get_shape_prst(shape) != 'roundRect':
                return False
            w, h = cm(shape.width), cm(shape.height)
            if not (3.4 <= w <= 3.8 and 1.3 <= h <= 1.6):
                return False
            # 位于界面相容优化模块内（左5–14.5cm，上10–17cm 大致范围）
            if not (5.0 <= cm(shape.left) and shape_right(shape) <= 14.5
                    and 10.0 <= cm(shape.top) and shape_bottom(shape) <= 17.0):
                return False
            lc = office_line_color(shape)
            lw = office_line_width_pt(shape)
            fc = office_fill(shape)
            return (color_is_blue_line(lc)
                    and lw is not None and 0.5 <= lw <= 1.0
                    and color_is_light_blue_or_white(fc))

        def shapes_overlap(rect, text_shape, tol=0.3):
            rl, rt = cm(rect.left), cm(rect.top)
            rr, rb = shape_right(rect), shape_bottom(rect)
            tl, tt = cm(text_shape.left), cm(text_shape.top)
            tr, tb = shape_right(text_shape), shape_bottom(text_shape)
            return not (rr < tl - tol or tr < rl - tol or rb < tt - tol or tb < rt - tol)

        rects = [s for s in slide1_shapes if small_roundrect_ok(s)]
        if len(rects) < 4:
            return False

        found_texts = set()
        used_rects = set()
        for txt in texts:
            for ts in slide1_shapes:
                t_content = get_text(ts)
                if txt not in t_content:
                    continue
                if not text_style_ok(ts):
                    continue
                for rect in rects:
                    if id(rect) in used_rects:
                        continue
                    if shapes_overlap(rect, ts):
                        found_texts.add(txt)
                        used_rects.add(id(rect))
                        break
                if txt in found_texts:
                    break

        if len(found_texts) < 4:
            return False

        # 2×2 排列验证：4个匹配矩形应有大致2列2行
        matched_rects = [s for s in rects if id(s) in used_rects]
        tops = sorted(set(round(cm(s.top), 1) for s in matched_rects))
        lefts = sorted(set(round(cm(s.left), 1) for s in matched_rects))
        return len(tops) >= 2 and len(lefts) >= 2

    score, desc = (3, "+3: 界面相容优化4个子项") if check_s1_interface_subitems() else (0, "+3: 界面相容优化4个子项")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +5: 界面相容优化子项左侧图标
    def check_s1_interface_icons():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(s):
            return cm(s.left) + cm(s.width)

        def shape_bottom(s):
            return cm(s.top) + cm(s.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
                'lt2': 'E8EFF7', 'bg2': 'E8EFF7',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def size_ok(shape):
            return 0.5 <= cm(shape.width) <= 0.8 and 0.5 <= cm(shape.height) <= 1.2

        def is_solid_blue(shape):
            fc = office_fill(shape)
            if not fc:
                return False
            try:
                r, g, b = int(fc[0:2], 16), int(fc[2:4], 16), int(fc[4:6], 16)
                return (color_is_blue(fc) or color_is_dark_blue(fc)) and not (b >= r and b >= g and (r + g + b) >= 580)
            except:
                return False

        def is_white_or_light(shape):
            fc = office_fill(shape)
            if not fc:
                return True
            try:
                r, g, b = int(fc[0:2], 16), int(fc[2:4], 16), int(fc[4:6], 16)
                return (r + g + b) >= 600
            except:
                return False

        # 根据4个文本位置，推断每个图标的约束区域（图标在文本左侧）
        # 以文本位置为锚点，不硬编码固定坐标
        def find_text_anchor(txt):
            for s in slide1_shapes:
                if txt in get_text(s) and 5.0 <= cm(s.left) <= 14.0 and 10.0 <= cm(s.top) <= 17.0:
                    return s
            return None

        def icon_region(text_shape, x_tol=1.5, y_tol=0.5):
            if text_shape is None:
                return None
            tl, tt = cm(text_shape.left), cm(text_shape.top)
            tb = shape_bottom(text_shape)
            return (tl - x_tol, tt - y_tol, tl + 0.1, tb + y_tol)

        t1 = find_text_anchor("增强表面亲水性")
        t2 = find_text_anchor("改进孔结构")
        t3 = find_text_anchor("膜基材")
        t4 = find_text_anchor("表面涂层")

        def in_region(shape, region):
            if region is None:
                return True  # 文本未定位时宽松
            lx, ly, rx, ry = region
            cx = cm(shape.left) + cm(shape.width) / 2
            cy = cm(shape.top) + cm(shape.height) / 2
            return lx <= cx <= rx + 0.5 and ly <= cy <= ry

        # ---- 1) 增强表面亲水性左侧：实心蓝色水滴（teardrop） ----
        r1 = icon_region(t1)
        droplet = any(
            get_shape_prst(s) == 'teardrop'
            and size_ok(s)
            and is_solid_blue(s)
            and in_region(s, r1)
            for s in slide1_shapes
        )

        # ---- 2) 改进孔结构左侧：大圆（外圆）+ 多个等大小蓝色小圆 ----
        r2 = icon_region(t2)
        pore_outer = [s for s in slide1_shapes
                      if get_shape_prst(s) == 'ellipse'
                      and in_region(s, r2)
                      and 0.4 <= cm(s.width) <= 0.8 and 0.4 <= cm(s.height) <= 0.8]
        pore_icon = False
        for outer in pore_outer:
            ol, ot = cm(outer.left), cm(outer.top)
            ow, oh = cm(outer.width), cm(outer.height)
            inner = [s for s in slide1_shapes
                     if get_shape_prst(s) == 'ellipse' and s != outer
                     and ol <= cm(s.left) and shape_right(s) <= ol + ow + 0.15
                     and ot <= cm(s.top) and shape_bottom(s) <= ot + oh + 0.15
                     and 0.05 <= cm(s.width) <= 0.25 and 0.05 <= cm(s.height) <= 0.25
                     and is_solid_blue(s)]
            if len(inner) >= 4:
                pore_icon = True
                break

        # ---- 3) 膜基材左侧：上蓝下白两个扁平矩形，上下有间隔 ----
        r3 = icon_region(t3)
        flat_blue = sorted(
            [s for s in slide1_shapes
             if get_shape_prst(s) == 'rect'
             and in_region(s, r3)
             and 0.4 <= cm(s.width) <= 0.85 and 0.05 <= cm(s.height) <= 0.35
             and is_solid_blue(s)],
            key=lambda s: cm(s.top)
        )
        flat_white = sorted(
            [s for s in slide1_shapes
             if get_shape_prst(s) == 'rect'
             and in_region(s, r3)
             and 0.4 <= cm(s.width) <= 0.85 and 0.05 <= cm(s.height) <= 0.35
             and is_white_or_light(s)],
            key=lambda s: cm(s.top)
        )
        layer_icon = False
        for bu in flat_blue:
            for wl in flat_white:
                gap = cm(wl.top) - shape_bottom(bu)
                if 0.0 < gap <= 0.3 and cm(bu.top) < cm(wl.top):
                    layer_icon = True
                    break
            if layer_icon:
                break

        # ---- 4) 表面涂层左侧：长方形（白色）+ 平行四边形（蓝色）拼接 ----
        r4 = icon_region(t4)
        coating_rects = [s for s in slide1_shapes
                         if get_shape_prst(s) == 'rect'
                         and in_region(s, r4)
                         and 0.4 <= cm(s.width) <= 0.85 and 0.1 <= cm(s.height) <= 0.65
                         and is_white_or_light(s)]
        coating_paras = [s for s in slide1_shapes
                         if get_shape_prst(s) == 'parallelogram'
                         and in_region(s, r4)
                         and 0.4 <= cm(s.width) <= 0.85 and 0.1 <= cm(s.height) <= 0.65
                         and is_solid_blue(s)]
        coating_icon = False
        for cr in coating_rects:
            for cp in coating_paras:
                gap = abs(shape_bottom(cr) - cm(cp.top))
                if gap <= 0.25 and cm(cr.top) < cm(cp.top):
                    coating_icon = True
                    break
            if coating_icon:
                break

        return droplet and pore_icon and layer_icon and coating_icon

    score, desc = (5, "+5: 界面相容优化图标") if check_s1_interface_icons() else (0, "+5: 界面相容优化图标")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: "运行抑污控制"模块
    def check_s1_operation_module():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
                'lt2': 'E8EFF7', 'bg2': 'E8EFF7',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def office_line_width_pt(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            if sp_pr is not None:
                ln = sp_pr.find(f'{DML}ln')
                if ln is not None:
                    w = ln.get('w')
                    if w:
                        return int(w) / 12700
            return get_line_width_pt(shape)

        def color_is_white_fill(c):
            if not c:
                return True  # 无填充色视为白色
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r >= 235 and g >= 235 and b >= 235
            except:
                return False

        def rpr_of(run):
            rpr = getattr(run._r, 'rPr', None)
            return rpr if rpr is not None else run._r.find(f'{DML}rPr')

        def def_rpr_of(para):
            p_pr = getattr(para._p, 'pPr', None)
            if p_pr is None:
                p_pr = para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def font_names_of(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    el = rpr.find(f'{DML}{tag}')
                    if el is not None and el.get('typeface'):
                        names.add(el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(n in ('Microsoft YaHei', '微软雅黑', 'SimHei', '黑体') for n in names)

        def size_pt_of(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def is_bold(run, def_rpr):
            try:
                if run.font.bold is True:
                    return True
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('b') in ('1', 'true'):
                    return True
            return False

        def run_color(run, def_rpr):
            try:
                if run.font.color.rgb:
                    return str(run.font.color.rgb)
            except:
                pass
            return color_from_elem(rpr_of(run)) or color_from_elem(def_rpr)

        def title_style_ok(shape):
            if get_text(shape) != "运行抑污控制":
                return False
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    def_rpr = def_rpr_of(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        size_pt = size_pt_of(run, def_rpr)
                        if not (font_name_ok(font_names_of(run, def_rpr))
                                and size_pt is not None and 14 <= size_pt <= 16
                                and is_bold(run, def_rpr)
                                and (color_is_blue(run_color(run, def_rpr))
                                     or color_is_dark_blue(run_color(run, def_rpr)))):
                            return False
                return has_run
            except:
                return False

        for box in slide1_shapes:
            if get_shape_prst(box) != 'roundRect':
                continue
            l, t, w, h = cm(box.left), cm(box.top), cm(box.width), cm(box.height)
            if not (13.0 <= l and shape_right(box) <= 24.0
                    and 10.0 <= t and shape_bottom(box) <= 16.0
                    and 7.0 <= w <= 8.0 and 4.0 <= h <= 5.2):
                continue
            lc = office_line_color(box)
            lw = office_line_width_pt(box)
            fc = office_fill(box)
            if not ((color_is_blue(lc) or color_is_dark_blue(lc))
                    and get_line_dash(box) == 'solid'
                    and lw is not None and 0.75 <= lw <= 1.25
                    and color_is_white_fill(fc)):
                continue
            for title in find_shapes_with_text(slide1_shapes, "运行抑污控制", partial=False):
                tl, tt = cm(title.left), cm(title.top)
                tr, tb = shape_right(title), shape_bottom(title)
                if l <= tl and tr <= l + w and t <= tt and tb <= t + h and title_style_ok(title):
                    return True
        return False

    score, desc = (1, "+1: 运行抑污控制模块") if check_s1_operation_module() else (0, "+1: 运行抑污控制模块")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: 运行抑污控制下4个小圆角矩形
    def check_s1_operation_subitems():
        texts = ["抑制浓差极化", "降低不可逆污染", "气水反冲", "运行参数"]
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
                'lt2': 'E8EFF7', 'bg2': 'E8EFF7',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def office_line_width_pt(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            if sp_pr is not None:
                ln = sp_pr.find(f'{DML}ln')
                if ln is not None:
                    w = ln.get('w')
                    if w:
                        return int(w) / 12700
            return get_line_width_pt(shape)

        def color_is_light_blue_or_white(c):
            if not c:
                return True
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                is_white = r >= 235 and g >= 235 and b >= 235
                is_light_blue = b >= r and b >= g and b >= 170 and (r + g + b) >= 500
                return is_white or is_light_blue
            except:
                return False

        def color_is_light_blue_line(c):
            return color_is_light_blue(c) or color_is_blue(c) or color_is_dark_blue(c)

        def rpr_of(run):
            rpr = getattr(run._r, 'rPr', None)
            return rpr if rpr is not None else run._r.find(f'{DML}rPr')

        def def_rpr_of(para):
            p_pr = getattr(para._p, 'pPr', None)
            if p_pr is None:
                p_pr = para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def font_names_of(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    el = rpr.find(f'{DML}{tag}')
                    if el is not None and el.get('typeface'):
                        names.add(el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(n in ('Microsoft YaHei', '微软雅黑', 'SimHei', '黑体') for n in names)

        def size_pt_of(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def item_style_ok(shape):
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    def_rpr = def_rpr_of(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        sz = size_pt_of(run, def_rpr)
                        if not (font_name_ok(font_names_of(run, def_rpr))
                                and sz is not None and 10 <= sz <= 12):
                            return False
                return has_run
            except:
                return False

        def small_roundrect_ok(box):
            if get_shape_prst(box) != 'roundRect':
                return False
            w, h = cm(box.width), cm(box.height)
            if not (3.4 <= w <= 3.8 and 1.3 <= h <= 1.6):
                return False
            if not (13.0 <= cm(box.left) and shape_right(box) <= 24.0
                    and 10.0 <= cm(box.top) and shape_bottom(box) <= 17.0):
                return False
            lc = office_line_color(box)
            lw = office_line_width_pt(box)
            fc = office_fill(box)
            return (color_is_light_blue_line(lc)
                    and lw is not None and 0.5 <= lw <= 1.0
                    and color_is_light_blue_or_white(fc))

        def shapes_overlap(rect, text_shape, tol=0.3):
            rl, rt, rr, rb = cm(rect.left), cm(rect.top), shape_right(rect), shape_bottom(rect)
            tl, tt, tr, tb = cm(text_shape.left), cm(text_shape.top), shape_right(text_shape), shape_bottom(text_shape)
            return not (rr < tl - tol or tr < rl - tol or rb < tt - tol or tb < rt - tol)

        rects = [s for s in slide1_shapes if small_roundrect_ok(s)]
        if len(rects) < 4:
            return False

        found_texts = set()
        used_rects = set()
        for txt in texts:
            for ts in slide1_shapes:
                if txt not in get_text(ts):
                    continue
                if not item_style_ok(ts):
                    continue
                for rect in rects:
                    if id(rect) in used_rects:
                        continue
                    if shapes_overlap(rect, ts):
                        found_texts.add(txt)
                        used_rects.add(id(rect))
                        break
                if txt in found_texts:
                    break

        if len(found_texts) < 4:
            return False

        matched_rects = [s for s in rects if id(s) in used_rects]
        tops = sorted(set(round(cm(s.top), 1) for s in matched_rects))
        lefts = sorted(set(round(cm(s.left), 1) for s in matched_rects))
        return len(tops) >= 2 and len(lefts) >= 2

    score, desc = (3, "+3: 运行抑污控制4个子项") if check_s1_operation_subitems() else (0, "+3: 运行抑污控制4个子项")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +5: 运行抑污控制子项图标
    def check_s1_operation_icons():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
                'lt2': 'E8EFF7', 'bg2': 'E8EFF7',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def icon_size_ok(shapes):
            left = min(cm(s.left) for s in shapes)
            top = min(cm(s.top) for s in shapes)
            right = max(cm(s.left) + cm(s.width) for s in shapes)
            bottom = max(cm(s.top) + cm(s.height) for s in shapes)
            return 0.5 <= right - left <= 0.8 and 0.6 <= bottom - top <= 0.9

        def is_solid_blue(shape):
            fc = office_fill(shape)
            return (color_is_dark_blue(fc) or color_is_blue(fc)) if fc else False

        def is_white_or_light_fill(shape):
            fc = office_fill(shape)
            if not fc:
                return True
            try:
                r, g, b = int(fc[0:2], 16), int(fc[2:4], 16), int(fc[4:6], 16)
                return (r + g + b) >= 600
            except:
                return False

        def is_blue_color(c):
            return color_is_blue(c) or color_is_dark_blue(c) or color_is_light_blue(c)

        def find_text_anchor(txt):
            for s in slide1_shapes:
                if txt in get_text(s) and 13.0 <= cm(s.left) <= 24.0 and 10.0 <= cm(s.top) <= 17.0:
                    return s
            return None

        def icon_region(text_shape, x_tol=1.5, y_tol=0.5):
            if text_shape is None:
                return None
            tl, tt = cm(text_shape.left), cm(text_shape.top)
            tb = cm(text_shape.top) + cm(text_shape.height)
            return (tl - x_tol, tt - y_tol, tl + 0.1, tb + y_tol)

        def in_region(shape, region):
            if region is None:
                return True
            lx, ly, rx, ry = region
            cx = cm(shape.left) + cm(shape.width) / 2
            cy = cm(shape.top) + cm(shape.height) / 2
            return lx <= cx <= rx + 0.5 and ly <= cy <= ry

        t1 = find_text_anchor("抑制浓差极化")
        t2 = find_text_anchor("降低不可逆污染")
        t3 = find_text_anchor("气水反冲")
        t4 = find_text_anchor("运行参数")

        def _center(shape):
            return (cm(shape.left) + cm(shape.width) / 2,
                    cm(shape.top) + cm(shape.height) / 2)

        def _bbox(shape):
            return (cm(shape.left), cm(shape.top),
                    cm(shape.left) + cm(shape.width),
                    cm(shape.top) + cm(shape.height))

        def _line_orientation(shape):
            # 返回 'h' / 'v' / None; 兼容用 spPr xfrm 或直线连接符
            w, h = cm(shape.width), abs(cm(shape.height))
            if w <= 0 and h <= 0:
                return None
            if w >= h * 2.0:
                return 'h'
            if h >= w * 2.0:
                return 'v'
            return None

        # ---- 1) 抑制浓差极化：一簇深蓝色实心小圆颗粒(分散排布) ----
        # 分散排布检查:颗粒中心 x/y 标准差都不能太小(避免"排一条线"),
        # 且颗粒 X 跨度和 Y 跨度都需占整体图标 bounds 的一定比例。
        r1 = icon_region(t1)
        particles = [s for s in slide1_shapes
                     if get_shape_prst(s) == 'ellipse'
                     and 0.04 <= cm(s.width) <= 0.20 and 0.04 <= cm(s.height) <= 0.20
                     and in_region(s, r1) and is_solid_blue(s)]
        concentration_icon = False
        if len(particles) >= 6 and icon_size_ok(particles):
            centers = [_center(p) for p in particles]
            xs = [c[0] for c in centers]
            ys = [c[1] for c in centers]
            n = len(centers)
            mean_x = sum(xs) / n
            mean_y = sum(ys) / n
            std_x = (sum((x - mean_x) ** 2 for x in xs) / n) ** 0.5
            std_y = (sum((y - mean_y) ** 2 for y in ys) / n) ** 0.5
            # 二维分散:两方向标准差都 >= 0.05cm,且跨度矩形近似为"簇"(短边/长边 >= 0.4)
            span_x = max(xs) - min(xs)
            span_y = max(ys) - min(ys)
            if (std_x >= 0.05 and std_y >= 0.05
                    and max(span_x, span_y) > 0
                    and min(span_x, span_y) / max(span_x, span_y) >= 0.4):
                concentration_icon = True

        # ---- 2) 降低不可逆污染：盾牌(竖线分左右,左侧中间横线) + 铅笔轮廓 ----
        # 拓扑要求:
        #   a. 盾牌外轮廓: pentagon/shield/homePlate 或近方形自由形状,占图标主体
        #   b. 盾牌内部一条纵向线,位于外轮廓横向中心 ±20% 处
        #   c. 该纵向线左侧有一条横向短线,位于纵向线上下中间 ±25%
        #   d. 图标区域内另有一个铅笔轮廓 shape(line/freeform/parallelogram),
        #      与盾牌外轮廓不同,尺寸细长 (长边/短边 >= 2)
        r2 = icon_region(t2)
        region_shapes = [s for s in slide1_shapes if in_region(s, r2)]
        shield_icon = False
        shield_candidates = [s for s in region_shapes
                             if get_shape_prst(s) in ('pentagon', 'shield',
                                                       'homePlate', 'flowChartAlternateProcess',
                                                       'rect', 'roundRect', 'freeform')
                             and cm(s.width) >= 0.3 and cm(s.height) >= 0.4]
        for shield in shield_candidates:
            sl, st, sr, sb = _bbox(shield)
            sw, sh = sr - sl, sb - st
            if sw <= 0 or sh <= 0:
                continue
            v_center = sl + sw / 2
            # 内部纵向线:方向 vertical,X 中心接近 v_center,Y 跨度覆盖盾牌高度 40%+
            v_lines = [s for s in region_shapes
                       if s is not shield
                       and get_shape_prst(s) == 'line'
                       and _line_orientation(s) == 'v'
                       and abs(_center(s)[0] - v_center) <= sw * 0.20
                       and abs(cm(s.height)) >= sh * 0.4
                       and st - 0.1 <= cm(s.top) <= sb]
            if not v_lines:
                continue
            vline = v_lines[0]
            vl, vt, vr, vb = _bbox(vline)
            vcx = (vl + vr) / 2
            vmid_y = (vt + vb) / 2
            # 左侧横向短线:方向 horizontal,X 中心在 vcx 左侧且在盾牌左半区内,
            # Y 接近纵向线的中点
            h_lines_left = [s for s in region_shapes
                            if s is not shield
                            and get_shape_prst(s) == 'line'
                            and _line_orientation(s) == 'h'
                            and _center(s)[0] < vcx
                            and _center(s)[0] >= sl - 0.05
                            and abs(_center(s)[1] - vmid_y) <= sh * 0.25]
            if not h_lines_left:
                continue
            # 铅笔轮廓:另一个细长形状,不是盾牌外轮廓、不是竖线、不是那些横线
            used = {id(shield), id(vline)} | {id(h) for h in h_lines_left}
            pencil = None
            for s in region_shapes:
                if id(s) in used:
                    continue
                if get_shape_prst(s) not in ('line', 'freeform', 'parallelogram',
                                              'rect', 'rightTriangle', 'triangle'):
                    continue
                w, h = cm(s.width), abs(cm(s.height))
                if max(w, h) < 0.15 or min(w, h) <= 0:
                    continue
                if max(w, h) / max(min(w, h), 0.01) >= 2.0:
                    pencil = s
                    break
            if pencil is None:
                continue
            shield_icon = icon_size_ok([shield, vline, h_lines_left[0], pencil])
            if shield_icon:
                break

        # ---- 3) 气水反冲：下方横向波浪线 + 上方大小不一"向上发散"的气泡 ----
        # 拓扑要求:
        #   a. 波浪线区域: 至少 2 条 horizontal line,Y 集中在图标下半部
        #   b. 气泡区域: 至少 3 个 ellipse 位于波浪线上方(Y 更小)
        #   c. 气泡尺寸不一致(极差 / 均值 >= 0.3),体现"大小不一"
        #   d. "向上发散": Y 越小(越靠上)的气泡在 X 上分布越宽 或至少气泡整体
        #      Y 跨度覆盖图标上半部分
        r3 = icon_region(t3)
        region3 = [s for s in slide1_shapes if in_region(s, r3)]
        wave_lines = [s for s in region3
                      if get_shape_prst(s) == 'line'
                      and _line_orientation(s) == 'h'
                      and cm(s.width) >= 0.2
                      and is_blue_color(office_line_color(s))]
        bubbles = [s for s in region3
                   if get_shape_prst(s) == 'ellipse'
                   and 0.05 <= cm(s.width) <= 0.35 and 0.05 <= cm(s.height) <= 0.35]
        backwash_icon = False
        if len(wave_lines) >= 2 and len(bubbles) >= 3:
            backwash_parts = wave_lines + bubbles
            if icon_size_ok(backwash_parts):
                wave_top = min(_center(w)[1] for w in wave_lines)
                bubbles_above = [b for b in bubbles if _center(b)[1] < wave_top + 0.02]
                if len(bubbles_above) >= 3:
                    diams = [(cm(b.width) + cm(b.height)) / 2 for b in bubbles_above]
                    d_mean = sum(diams) / len(diams)
                    d_range = max(diams) - min(diams)
                    # "大小不一": 极差 / 均值 >= 0.3
                    if d_mean > 0 and d_range / d_mean >= 0.3:
                        # 向上发散:气泡整体 Y 跨度 >= 0.2cm(至少有多层向上排布)
                        by = [_center(b)[1] for b in bubbles_above]
                        if max(by) - min(by) >= 0.2:
                            backwash_icon = True

        # ---- 4) 运行参数：上下两条平行水平短线 + 中间竖线串联 3 个空心小圆 ----
        # 拓扑要求:
        #   a. 恰好 2 条 horizontal line, Y 有明显上下之分, 长度相近(极差 <=40% 均值)
        #   b. >=1 条 vertical line, 位于两条 h line 之间的中段位置
        #   c. 恰好 3 个空心/白色 ellipse, X 中心与竖线接近(±平均圆直径),
        #      Y 分布在两条 h line 之间
        r4 = icon_region(t4)
        region4 = [s for s in slide1_shapes if in_region(s, r4)]
        param_h_lines = [s for s in region4
                         if get_shape_prst(s) == 'line'
                         and _line_orientation(s) == 'h'
                         and cm(s.width) >= 0.2]
        param_v_lines = [s for s in region4
                         if get_shape_prst(s) == 'line'
                         and _line_orientation(s) == 'v']
        param_circles = [s for s in region4
                         if get_shape_prst(s) == 'ellipse'
                         and 0.05 <= cm(s.width) <= 0.25 and 0.05 <= cm(s.height) <= 0.25
                         and is_white_or_light_fill(s)]
        param_icon = False
        if len(param_h_lines) >= 2 and len(param_v_lines) >= 1 and len(param_circles) == 3:
            # 取最像上下两条平行的两条 h line: 按 Y 排序后取首尾
            h_sorted = sorted(param_h_lines, key=lambda s: _center(s)[1])
            top_h, bot_h = h_sorted[0], h_sorted[-1]
            top_y = _center(top_h)[1]
            bot_y = _center(bot_h)[1]
            if bot_y - top_y >= 0.15:  # 上下必须有明显间距
                widths = [cm(top_h.width), cm(bot_h.width)]
                w_mean = sum(widths) / 2
                if w_mean > 0 and abs(widths[0] - widths[1]) <= w_mean * 0.4:
                    mid_y = (top_y + bot_y) / 2
                    # 竖线应位于两条 h line 之间(其 Y 中心在 [top_y, bot_y] 内 ±10%)
                    inner_v = [v for v in param_v_lines
                               if top_y - 0.05 <= _center(v)[1] <= bot_y + 0.05]
                    if inner_v:
                        vcx = _center(inner_v[0])[0]
                        avg_diam = sum(cm(c.width) for c in param_circles) / 3
                        # 3 个空心圆:X 中心与竖线 X 接近,Y 都在两条 h line 之间
                        aligned_circles = [c for c in param_circles
                                           if abs(_center(c)[0] - vcx) <= max(avg_diam, 0.15)
                                           and top_y - 0.05 <= _center(c)[1] <= bot_y + 0.05]
                        if len(aligned_circles) == 3:
                            cy_sorted = sorted(_center(c)[1] for c in aligned_circles)
                            gap1 = cy_sorted[1] - cy_sorted[0]
                            gap2 = cy_sorted[2] - cy_sorted[1]
                            g_mean = (gap1 + gap2) / 2
                            # 3 个圆纵向大致均匀(相邻间距差 <= 50% 均值),体现"串联"
                            if g_mean > 0 and abs(gap1 - gap2) <= g_mean * 0.5:
                                parts = [top_h, bot_h, inner_v[0]] + aligned_circles
                                param_icon = icon_size_ok(parts)

        return concentration_icon and shield_icon and backwash_icon and param_icon

    score, desc = (5, "+5: 运行抑污控制图标") if check_s1_operation_icons() else (0, "+5: 运行抑污控制图标")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 两模块下方连接结构
    def check_s1_bottom_connector():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def shape_center_x(shape):
            return cm(shape.left) + cm(shape.width) / 2

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def blue(shape):
            return color_is_blue(office_line_color(shape)) or color_is_dark_blue(office_line_color(shape))

        def blue_fill_or_line(shape):
            return (color_is_blue(office_fill(shape)) or color_is_dark_blue(office_fill(shape))
                    or color_is_blue(office_line_color(shape)) or color_is_dark_blue(office_line_color(shape)))

        # 结果栏：用于确认箭头指向含有指定文本的长条圆角矩形
        result_text = "通量衰减减缓·清洗周期延长·分离效率提升"
        result_targets = [s for s in slide1_shapes if result_text in get_text(s)]
        if not result_targets:
            return False
        result_top = min(cm(s.top) for s in result_targets)
        result_center_x = sum(shape_center_x(s) for s in result_targets) / len(result_targets)

        blue_solid_lines = [s for s in slide1_shapes
                            if get_shape_prst(s) == 'line'
                            and get_line_dash(s) in ('solid', None)
                            and blue(s)]

        # 两模块下方各一根竖向高0.4-0.6cm蓝色直实线
        verticals = [s for s in blue_solid_lines
                     if abs(cm(s.width)) <= 0.2
                     and 0.4 <= abs(cm(s.height)) <= 0.6
                     and 15.5 <= cm(s.top) <= 18.0]
        if len(verticals) < 2:
            return False

        # 一根长6-8cm蓝色实心横向直线
        horizontals = [s for s in blue_solid_lines
                       if 6.0 <= abs(cm(s.width)) <= 8.0
                       and abs(cm(s.height)) <= 0.2
                       and 16.0 <= cm(s.top) <= 18.0]
        if not horizontals:
            return False

        for hline in horizontals:
            h_left = min(cm(hline.left), shape_right(hline))
            h_right = max(cm(hline.left), shape_right(hline))
            h_y = cm(hline.top)
            left_verticals = [v for v in verticals
                              if abs(shape_center_x(v) - h_left) <= 0.6
                              and abs(shape_bottom(v) - h_y) <= 0.6]
            right_verticals = [v for v in verticals
                               if abs(shape_center_x(v) - h_right) <= 0.6
                               and abs(shape_bottom(v) - h_y) <= 0.6]
            if not (left_verticals and right_verticals):
                continue

            # 横线中心二分之一处有竖向蓝色实心箭头，高0.2-0.6cm，指向结果栏
            h_mid = (h_left + h_right) / 2
            arrows = []
            for s in slide1_shapes:
                prst = get_shape_prst(s)
                if prst == 'downArrow':
                    if not (0.2 <= cm(s.height) <= 0.6 and blue_fill_or_line(s)):
                        continue
                    sx = shape_center_x(s)
                    sy = shape_bottom(s)
                elif is_line_or_connector(s):
                    if not (abs(cm(s.width)) <= 0.2 and 0.2 <= abs(cm(s.height)) <= 0.6 and blue(s)):
                        continue
                    if not has_line_arrowhead(s):
                        continue
                    sx = shape_center_x(s)
                    sy = shape_bottom(s)
                else:
                    continue
                if abs(sx - h_mid) <= 0.6 and h_y - 0.2 <= cm(s.top) <= h_y + 0.8 and sy <= result_top + 0.5:
                    arrows.append(s)
            if arrows and abs(h_mid - result_center_x) <= 2.0:
                return True
        return False

    score, desc = (1, "+1: 两模块下方连接线") if check_s1_bottom_connector() else (0, "+1: 两模块下方连接线")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 底部长条结果栏
    def check_s1_result_bar():
        expected_text = "通量衰减减缓·清洗周期延长·分离效率提升"
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '1F4E79', 'accent1': '4472C4', 'accent5': '5B9BD5',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
                'lt2': 'E8EFF7', 'bg2': 'E8EFF7',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def office_line_width_pt(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            if sp_pr is not None:
                ln = sp_pr.find(f'{DML}ln')
                if ln is not None:
                    w = ln.get('w')
                    if w:
                        return int(w) / 12700
            return get_line_width_pt(shape)

        def color_is_light_blue_fill(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return b >= r and b >= g and b >= 170 and (r + g + b) >= 500
            except:
                return False

        def rpr_of(run):
            rpr = getattr(run._r, 'rPr', None)
            return rpr if rpr is not None else run._r.find(f'{DML}rPr')

        def def_rpr_of(para):
            p_pr = getattr(para._p, 'pPr', None)
            if p_pr is None:
                p_pr = para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def font_names_of(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    el = rpr.find(f'{DML}{tag}')
                    if el is not None and el.get('typeface'):
                        names.add(el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(n in ('Microsoft YaHei', '微软雅黑', 'SimHei', '黑体') for n in names)

        def size_pt_of(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def run_color(run, def_rpr):
            try:
                if run.font.color.rgb:
                    return str(run.font.color.rgb)
            except:
                pass
            return color_from_elem(rpr_of(run)) or color_from_elem(def_rpr)

        def text_style_ok(shape):
            if get_text(shape) != expected_text:
                return False
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    def_rpr = def_rpr_of(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        size_pt = size_pt_of(run, def_rpr)
                        if not (font_name_ok(font_names_of(run, def_rpr))
                                and size_pt is not None and 13 <= size_pt <= 15
                                and color_is_dark_blue(run_color(run, def_rpr))):
                            return False
                return has_run
            except:
                return False

        for box in slide1_shapes:
            if get_shape_prst(box) != 'roundRect':
                continue
            l, t, w, h = cm(box.left), cm(box.top), cm(box.width), cm(box.height)
            if not (7.0 <= l and shape_right(box) <= 22.0
                    and 15.0 <= t and shape_bottom(box) <= 18.0
                    and 12.0 <= w <= 14.0 and 0.8 <= h <= 1.2):
                continue
            lc = office_line_color(box)
            lw = office_line_width_pt(box)
            fc = office_fill(box)
            if not ((color_is_blue(lc) or color_is_dark_blue(lc))
                    and get_line_dash(box) == 'solid'
                    and lw is not None and 0.75 <= lw <= 1.25
                    and color_is_light_blue_fill(fc)):
                continue
            for ts in find_shapes_with_text(slide1_shapes, expected_text, partial=False):
                tl, tt = cm(ts.left), cm(ts.top)
                tr, tb = shape_right(ts), shape_bottom(ts)
                if l <= tl and tr <= l + w and t <= tt and tb <= t + h and text_style_ok(ts):
                    return True
        return False

    score, desc = (1, "+1: 底部结果栏") if check_s1_result_bar() else (0, "+1: 底部结果栏")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +5: 右侧膜分离设备
    def check_s1_device():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(s):
            return cm(s.left) + cm(s.width)

        def shape_bottom(s):
            return cm(s.top) + cm(s.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'dk2': '404040', 'accent3': 'A5A5A5', 'accent4': 'BFBFBF',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
                'lt2': 'D9D9D9', 'bg2': 'D9D9D9',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def office_line_color(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            ln = sp_pr.find(f'{DML}ln') if sp_pr is not None else None
            return color_from_elem(ln) or get_line_color(shape)

        def office_line_width_pt(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            if sp_pr is not None:
                ln = sp_pr.find(f'{DML}ln')
                if ln is not None:
                    w = ln.get('w')
                    if w:
                        return int(w) / 12700
            return get_line_width_pt(shape)

        def color_is_gray_or_light(c):
            if not c:
                return True
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                is_gray = abs(r - g) <= 30 and abs(g - b) <= 30 and abs(r - b) <= 30
                is_light = (r + g + b) >= 450
                return is_gray or is_light
            except:
                return False

        def color_is_dark_gray(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return abs(r - g) <= 40 and abs(g - b) <= 40 and max(r, g, b) <= 130
            except:
                return False

        def color_is_white_or_near_white(c):
            if not c:
                return True
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r >= 230 and g >= 220 and b >= 200
            except:
                return False

        def lw_ok(shape):
            lw = office_line_width_pt(shape)
            return lw is not None and 0.5 <= lw <= 1.25

        def in_device_area(shape):
            return (21.0 <= cm(shape.left) and shape_right(shape) <= 27.0
                    and 4.5 <= cm(shape.top) and shape_bottom(shape) <= 17.0)

        region = [s for s in slide1_shapes if in_device_area(s)]
        if len(region) < 10:
            return False

        # 整体尺寸约束
        overall_left = min(cm(s.left) for s in region)
        overall_top = min(cm(s.top) for s in region)
        overall_right = max(shape_right(s) for s in region)
        overall_bottom = max(shape_bottom(s) for s in region)
        if not (21.0 <= overall_left and overall_right <= 27.0
                and 4.5 <= overall_top and overall_bottom <= 17.0
                and 2.3 <= overall_right - overall_left <= 3.0
                and 10.0 <= overall_bottom - overall_top <= 13.0):
            return False

        # 1) 最上方圆形压力表表盘（浅灰色圆形，宽高0.4-1.0cm）
        gauge = any(
            get_shape_prst(s) in ('ellipse', 'arc')
            and 0.4 <= cm(s.width) <= 1.0 and 0.4 <= cm(s.height) <= 1.0
            and cm(s.top) <= overall_top + 2.0
            and color_is_gray_or_light(office_fill(s))
            and lw_ok(s)
            for s in region
        )
        if not gauge:
            return False

        # 2) 主体空心长圆筒（高宽比大，宽2-2.8cm，高5cm以上，灰色填充，深灰外轮廓）
        body = any(
            get_shape_prst(s) in ('rect', 'roundRect', 'ellipse')
            and 2.0 <= cm(s.width) <= 2.8
            and cm(s.height) >= 5.0
            and color_is_gray_or_light(office_fill(s))
            and color_is_dark_gray(office_line_color(s))
            and lw_ok(s)
            for s in region
        )
        if not body:
            return False

        # 3) 筒体内密集等距竖向细长白色竖线（白色/浅米色，宽极细，高度较长）
        inner_lines = [s for s in region
                       if get_shape_prst(s) in ('rect', 'line')
                       and cm(s.width) <= 0.15
                       and cm(s.height) >= 2.0
                       and color_is_white_or_near_white(office_fill(s) or office_line_color(s))]
        if len(inner_lines) < 5:
            return False

        # 4) 左侧侧管（水平短圆柱管道），椭圆/圆形管口，宽高约0.3-0.8cm
        side_ports = [s for s in region
                      if get_shape_prst(s) in ('ellipse', 'rect', 'roundRect')
                      and 0.2 <= cm(s.width) <= 0.9 and 0.2 <= cm(s.height) <= 0.9
                      and cm(s.left) <= overall_left + 0.8
                      and overall_top + 1.5 <= cm(s.top) <= overall_bottom - 1.5]
        if len(side_ports) < 2:
            return False

        # 5) 底部扁平圆柱体和突出半宽圆柱体（宽与筒体相近或为其一半，高较小）
        bottom_parts = [s for s in region
                        if get_shape_prst(s) in ('rect', 'roundRect', 'ellipse')
                        and 0.8 <= cm(s.width) <= 2.8
                        and 0.15 <= cm(s.height) <= 1.5
                        and cm(s.top) >= overall_bottom - 3.0]
        if len(bottom_parts) < 2:
            return False

        # 6) 底部三根立式圆柱支脚（细长矩形或线，位于底部）
        legs = [s for s in region
                if get_shape_prst(s) in ('rect', 'line')
                and cm(s.width) <= 0.5
                and cm(s.height) >= 0.3
                and cm(s.top) >= overall_bottom - 2.0]
        return len(legs) >= 3

    score, desc = (5, "+5: 右侧膜分离设备图") if check_s1_device() else (0, "+5: 右侧膜分离设备图")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 右侧红色竖排文本
    def check_s1_red_text():
        expected_chars = "降低膜污染并提升分离效率"
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'accent2': 'ED7D31', 'accent4': 'FF0000',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def rpr_of(run):
            rpr = getattr(run._r, 'rPr', None)
            return rpr if rpr is not None else run._r.find(f'{DML}rPr')

        def def_rpr_of(para):
            p_pr = getattr(para._p, 'pPr', None)
            if p_pr is None:
                p_pr = para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def font_names_of(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    el = rpr.find(f'{DML}{tag}')
                    if el is not None and el.get('typeface'):
                        names.add(el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(n in ('Microsoft YaHei', '微软雅黑', 'SimHei', '黑体') for n in names)

        def size_pt_of(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def is_bold(run, def_rpr):
            try:
                if run.font.bold is True:
                    return True
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('b') in ('1', 'true'):
                    return True
            return False

        def run_color(run, def_rpr):
            try:
                if run.font.color.rgb:
                    return str(run.font.color.rgb)
            except:
                pass
            return color_from_elem(rpr_of(run)) or color_from_elem(def_rpr)

        def color_is_red_or_dark_red(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r >= 150 and g <= 80 and b <= 80 and r > g and r > b
            except:
                return False

        def text_has_vertical_direction(shape):
            try:
                txBody = shape._element.find(f'.//{DML}txBody')
                if txBody is None:
                    return False
                bodyPr = txBody.find(f'{DML}bodyPr')
                if bodyPr is not None:
                    vert = bodyPr.get('vert', '')
                    if vert in ('vert', 'vert270', 'eaVert', 'mongolianVert'):
                        return True
            except:
                pass
            return False

        def text_is_single_column_vertical(shape):
            txt = get_text(shape)
            lines = [l for l in shape.text.split('\n') if l.strip()] if hasattr(shape, 'text') else []
            if len(lines) >= 4 and all(len(l.strip()) == 1 for l in lines):
                return True
            return False

        def in_position(shape):
            return (26.0 <= cm(shape.left) and shape_right(shape) <= 28.0
                    and 4.0 <= cm(shape.top) and shape_bottom(shape) <= 15.0)

        def style_ok(shape):
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    def_rpr = def_rpr_of(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        sz = size_pt_of(run, def_rpr)
                        if not (font_name_ok(font_names_of(run, def_rpr))
                                and sz is not None and 18 <= sz <= 20
                                and is_bold(run, def_rpr)
                                and color_is_red_or_dark_red(run_color(run, def_rpr))):
                            return False
                return has_run
            except:
                return False

        for shape in slide1_shapes:
            if not in_position(shape):
                continue
            txt = get_text(shape)
            # 竖排文本框（含全部文字）或逐字竖排（每行一字）
            is_vertical = (text_has_vertical_direction(shape)
                           or text_is_single_column_vertical(shape)
                           or (hasattr(shape, 'text') and shape.text
                               and all(c in shape.text for c in expected_chars)
                               and cm(shape.width) < cm(shape.height)))
            if not is_vertical:
                continue
            has_all_chars = all(c in txt for c in expected_chars)
            if not has_all_chars:
                continue
            if style_ok(shape):
                return True

        # 逐字一个形状排列方式：多个独立单字形状叠加
        char_shapes = []
        for shape in slide1_shapes:
            if not in_position(shape):
                continue
            txt = get_text(shape)
            if len(txt) == 1 and txt in expected_chars and style_ok(shape):
                char_shapes.append(txt)
        found_chars = set(char_shapes)
        return len(found_chars) >= len(expected_chars) * 0.8


    score, desc = (1, "+1: 右侧红色竖排文本") if check_s1_red_text() else (0, "+1: 右侧红色竖排文本")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 底部图题
    def check_s1_title():
        expected_title = "图2 膜分离水处理系统污染控制与通量提升研究"
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def shape_center_x(shape):
            return cm(shape.left) + cm(shape.width) / 2

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def rpr_of(run):
            rpr = getattr(run._r, 'rPr', None)
            return rpr if rpr is not None else run._r.find(f'{DML}rPr')

        def def_rpr_of(para):
            p_pr = getattr(para._p, 'pPr', None)
            if p_pr is None:
                p_pr = para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def font_names_of(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    el = rpr.find(f'{DML}{tag}')
                    if el is not None and el.get('typeface'):
                        names.add(el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(n in ('SimSun', '宋体', 'SimHei', '黑体', 'Microsoft YaHei', '微软雅黑') for n in names)

        def size_pt_of(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def is_bold_or_semibold(run, def_rpr):
            try:
                if run.font.bold is True:
                    return True
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is None:
                    continue
                if rpr.get('b') in ('1', 'true'):
                    return True
                # Office 可用字体粗细 w=600/700 表示半粗/粗
                if rpr.get('w') and int(rpr.get('w')) >= 600:
                    return True
            return False

        def run_color(run, def_rpr):
            try:
                if run.font.color.rgb:
                    return str(run.font.color.rgb)
            except:
                pass
            return color_from_elem(rpr_of(run)) or color_from_elem(def_rpr)

        def color_is_black(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return max(r, g, b) <= 35
            except:
                return False

        def para_centered(para):
            try:
                from pptx.enum.text import PP_ALIGN
                if para.alignment == PP_ALIGN.CENTER:
                    return True
            except:
                pass
            p_pr = para._p.find(f'{DML}pPr')
            return p_pr is not None and p_pr.get('algn') == 'ctr'

        def title_style_ok(shape):
            try:
                has_run = False
                centered = False
                for para in shape.text_frame.paragraphs:
                    if para_centered(para):
                        centered = True
                    def_rpr = def_rpr_of(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        size_pt = size_pt_of(run, def_rpr)
                        if not (font_name_ok(font_names_of(run, def_rpr))
                                and size_pt is not None and 15 <= size_pt <= 17
                                and color_is_black(run_color(run, def_rpr))
                                and is_bold_or_semibold(run, def_rpr)):
                            return False
                # 段落居中，或形状中心接近指定区域中心(15cm)
                horizontal_centered = centered or abs(shape_center_x(shape) - 15.0) <= 0.6
                return has_run and horizontal_centered
            except:
                return False

        for m in find_shapes_with_text(slide1_shapes, expected_title, partial=False):
            if not (8.0 <= cm(m.left) and shape_right(m) <= 22.0
                    and 18.0 <= cm(m.top) and shape_bottom(m) <= 20.0):
                continue
            if title_style_ok(m):
                return True
        return False

    score, desc = (1, "+1: 第1页底部图题") if check_s1_title() else (0, "+1: 第1页底部图题")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # ===== 第2页评分 =====

    # +5: 顶部4段流程箭头色块
    def check_s2_flow_blocks():
        texts = ["数据获取与治理", "特征提取与建模", "模型评估与解释", "干预建议生成"]
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def base_scheme_color(val):
            scheme_map = {
                'dk1': '000000', 'tx1': '000000',
                'accent1': '4472C4', 'accent2': 'ED7D31',
                'accent3': 'A9D18E', 'accent4': 'FF0000',
                'accent5': '5B9BD5', 'accent6': '70AD47',
                'lt1': 'FFFFFF', 'bg1': 'FFFFFF',
            }
            return scheme_map.get((val or '').lower()) or STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                val = int(child.get('val') or 100000) / 100000
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def color_from_elem(parent):
            if parent is None:
                return None
            solid = parent.find(f'{DML}solidFill')
            if solid is None:
                return None
            srgb = solid.find(f'{DML}srgbClr')
            if srgb is not None:
                return apply_transforms(srgb.get('val'), srgb)
            scheme = solid.find(f'{DML}schemeClr')
            if scheme is not None:
                return apply_transforms(base_scheme_color(scheme.get('val')), scheme)
            return None

        def office_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return color_from_elem(sp_pr) or get_fill_color(shape)

        def rpr_of(run):
            rpr = getattr(run._r, 'rPr', None)
            return rpr if rpr is not None else run._r.find(f'{DML}rPr')

        def def_rpr_of(para):
            p_pr = getattr(para._p, 'pPr', None)
            if p_pr is None:
                p_pr = para._p.find(f'{DML}pPr')
            return p_pr.find(f'{DML}defRPr') if p_pr is not None else None

        def font_names_of(run, def_rpr):
            names = set()
            try:
                if run.font.name:
                    names.add(run.font.name)
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is None:
                    continue
                for tag in ('latin', 'ea', 'cs'):
                    el = rpr.find(f'{DML}{tag}')
                    if el is not None and el.get('typeface'):
                        names.add(el.get('typeface'))
            return names

        def font_name_ok(names):
            return any(n in ('Microsoft YaHei', '微软雅黑', 'SimHei', '黑体') for n in names)

        def size_pt_of(run, def_rpr):
            try:
                if run.font.size:
                    return run.font.size / 12700
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('sz'):
                    return int(rpr.get('sz')) / 100
            return None

        def is_bold(run, def_rpr):
            try:
                if run.font.bold is True:
                    return True
            except:
                pass
            for rpr in (rpr_of(run), def_rpr):
                if rpr is not None and rpr.get('b') in ('1', 'true'):
                    return True
            return False

        def run_color(run, def_rpr):
            try:
                if run.font.color.rgb:
                    return str(run.font.color.rgb)
            except:
                pass
            return color_from_elem(rpr_of(run)) or color_from_elem(def_rpr)

        def color_is_white_text(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r >= 240 and g >= 240 and b >= 240
            except:
                return False

        def text_style_ok(shape):
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    def_rpr = def_rpr_of(para)
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        sz = size_pt_of(run, def_rpr)
                        if not (font_name_ok(font_names_of(run, def_rpr))
                                and sz is not None and 14 <= sz <= 16
                                and is_bold(run, def_rpr)
                                and color_is_white_text(run_color(run, def_rpr))):
                            return False
                return has_run
            except:
                return False

        # 色块形状：chevron（左圆右箭头）或 rightArrow/pentagon；也兼容普通矩形
        allowed_prsts = ('chevron', 'rightArrow', 'pentagon', 'flowChartProcess', 'rect', 'roundRect')
        blocks = sorted(
            [s for s in slide2_shapes
             if get_shape_prst(s) in allowed_prsts
             and 2.0 <= cm(s.left) and shape_right(s) <= 29.0
             and 0.8 <= cm(s.top) and shape_bottom(s) <= 2.8
             and 6.0 <= cm(s.width) <= 6.5
             and 1.4 <= cm(s.height) <= 1.7],
            key=lambda s: cm(s.left)
        )
        if len(blocks) < 4:
            return False

        # 整体位置约束
        overall_left = min(cm(s.left) for s in blocks[:4])
        overall_right = max(shape_right(s) for s in blocks[:4])
        overall_top = min(cm(s.top) for s in blocks[:4])
        overall_bottom = max(shape_bottom(s) for s in blocks[:4])
        if not (2.0 <= overall_left and overall_right <= 29.0
                and 0.8 <= overall_top and overall_bottom <= 2.8):
            return False

        expected_colors = [color_is_green, color_is_blue, color_is_purple, color_is_orange]

        for block, txt, color_check in zip(blocks[:4], texts, expected_colors):
            fc = office_fill(block)
            if not color_check(fc):
                return False
            bl, bt = cm(block.left), cm(block.top)
            br, bb = shape_right(block), shape_bottom(block)
            found_text = False
            for ts in slide2_shapes:
                if txt not in get_text(ts):
                    continue
                tl, tt, tr, tb = cm(ts.left), cm(ts.top), shape_right(ts), shape_bottom(ts)
                if bl <= tl and tr <= br + 0.3 and bt <= tt and tb <= bb + 0.3 and text_style_ok(ts):
                    found_text = True
                    break
            if not found_text:
                return False
        return True

    score, desc = (5, "+5: 第2页顶部4段流程箭头色块") if check_s2_flow_blocks() else (0, "+5: 第2页顶部4段流程箭头色块")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +5: 流程色块间箭头
    def check_s2_flow_arrows():
        # 三个连接箭头必须宽0.9-1.3cm，颜色依次绿/紫/黄，且为渐变填充。
        # rubric 强制"渐变式填充": 只认块状箭头预设 + spPr 内 gradFill,
        # 不再接受 line/连接符线色兜底,也不再接受 solidFill 兜底。
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def base_scheme_color(val):
            return STANDARD_SCHEME_COLORS.get((val or '').lower())

        def apply_transforms(color, element):
            if not color:
                return None
            r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
            for child in element:
                tag = child.tag.split('}')[-1]
                try:
                    val = int(child.get('val') or 100000) / 100000
                except (TypeError, ValueError):
                    continue
                if tag == 'lumMod':
                    r, g, b = r * val, g * val, b * val
                elif tag == 'lumOff':
                    r, g, b = r + 255 * val, g + 255 * val, b + 255 * val
                elif tag == 'tint':
                    r, g, b = r + (255 - r) * val, g + (255 - g) * val, b + (255 - b) * val
                elif tag == 'shade':
                    r, g, b = r * val, g * val, b * val
            r, g, b = [max(0, min(255, int(round(v)))) for v in (r, g, b)]
            return f'{r:02X}{g:02X}{b:02X}'

        def get_grad_stop_colors(shape):
            """从形状 spPr 层级的 gradFill 中提取所有 stop 的颜色 (支持 srgbClr / schemeClr)。"""
            sp = shape._element
            sp_pr = sp.find(f'.//{DML}spPr')
            if sp_pr is None:
                return []
            grad_fill = sp_pr.find(f'{DML}gradFill')
            if grad_fill is None:
                return []
            colors = []
            for stop in grad_fill.findall(f'.//{DML}gs'):
                srgb = stop.find(f'{DML}srgbClr')
                if srgb is not None:
                    colors.append(apply_transforms(srgb.get('val'), srgb))
                    continue
                scheme = stop.find(f'{DML}schemeClr')
                if scheme is not None:
                    colors.append(apply_transforms(base_scheme_color(scheme.get('val')),
                                                    scheme))
            return [c for c in colors if c]

        def has_grad_fill(shape):
            """必须是 spPr 层级 (即形状本身) 的 gradFill,而不是其他祖先。"""
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return sp_pr is not None and sp_pr.find(f'{DML}gradFill') is not None

        def grad_color_matches(shape, color_ok):
            return any(color_ok(c) for c in get_grad_stop_colors(shape))

        # 只认块状箭头预设,不再接受 line / 直线连接符
        arrow_prsts = ('rightArrow', 'chevron', 'bentArrow', 'stripedRightArrow',
                       'notchedRightArrow', 'homePlate', 'pentagon', 'flowChartProcess')

        # 三个箭头位置区间(left)及对应颜色判断函数
        expected = [
            (7.4, 8.8, color_is_green),
            (14.8, 16.2, color_is_purple),
            (22.0, 23.5, color_is_yellow),
        ]
        found = 0
        for left_lo, left_hi, color_ok in expected:
            ok = False
            for s in slide2_shapes:
                l = cm(s.left)
                t = cm(s.top)
                w = abs(cm(s.width))
                if not (left_lo <= l <= left_hi and 1.0 <= t <= 2.4 and 0.9 <= w <= 1.3):
                    continue
                prst = get_shape_prst(s)
                if prst not in arrow_prsts:
                    continue
                # 严格要求渐变填充 + 渐变 stop 颜色命中
                if not has_grad_fill(s):
                    continue
                if not grad_color_matches(s, color_ok):
                    continue
                ok = True
                break
            if ok:
                found += 1
        return found == 3

    score, desc = (5, "+5: 流程色块间渐变箭头") if check_s2_flow_arrows() else (0, "+5: 流程色块间渐变箭头")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +5: 流程色块左侧圆形图标
    def check_s2_flow_icons():
        # 四个流程块左侧必须有0.7-1.2cm的圆形白底图标区域，并包含题干指定的可编辑线性图标
        def iter_office_shapes(shapes):
            for shape in shapes:
                yield shape
                if hasattr(shape, 'shapes'):
                    for child in iter_office_shapes(shape.shapes):
                        yield child

        shapes = list(iter_office_shapes(slide2_shapes))

        def right(shape):
            return cm(shape.left) + cm(shape.width)

        def bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def center(shape):
            return cm(shape.left) + cm(shape.width) / 2, cm(shape.top) + cm(shape.height) / 2

        def size_close(a, b, tolerance=0.12):
            return abs(cm(a.width) - cm(b.width)) <= tolerance and abs(cm(a.height) - cm(b.height)) <= tolerance

        def in_circle(shape, circle, margin=0.05):
            cl, ct, cr, cb = cm(circle.left), cm(circle.top), right(circle), bottom(circle)
            return (cl - margin <= cm(shape.left) and right(shape) <= cr + margin
                    and ct - margin <= cm(shape.top) and bottom(shape) <= cb + margin)

        def line_len(shape):
            return (cm(shape.width) ** 2 + cm(shape.height) ** 2) ** 0.5

        def is_line(shape):
            return get_shape_prst(shape) == 'line' and not get_text(shape)

        def is_small_dot(shape):
            return (get_shape_prst(shape) == 'ellipse' and not get_text(shape)
                    and 0.08 <= cm(shape.width) <= 0.28 and 0.08 <= cm(shape.height) <= 0.28)

        def icon_parts(circle):
            return [s for s in shapes if s != circle and in_circle(s, circle) and not get_text(s)]

        def has_wireless_and_wave(circle):
            parts = icon_parts(circle)
            arcs = [s for s in parts if get_shape_prst(s) == 'arc']
            short_lines = [s for s in parts if is_line(s) and line_len(s) >= 0.08]
            dots = [s for s in parts if is_small_dot(s)]
            # 无线信号：办公软件可编辑弧线至少2条，或1条弧线+信号点；数据波形：多段可编辑折线/线段
            return len(arcs) >= 2 and len(short_lines) >= 3 and len(dots) >= 1

        def has_five_node_network(circle):
            parts = icon_parts(circle)
            dots = [s for s in parts if is_small_dot(s)]
            lines = [s for s in parts if is_line(s) and line_len(s) >= 0.18]
            if len(dots) < 6 or len(lines) < 5:
                return False
            cx, cy = center(circle)
            center_dots = [d for d in dots if abs(center(d)[0] - cx) <= 0.18 and abs(center(d)[1] - cy) <= 0.18]
            if not center_dots:
                return False
            center_dot = center_dots[0]
            same_size_dots = [d for d in dots if size_close(d, center_dot)]
            outer_dots = [d for d in same_size_dots if d != center_dot]
            # 中间一个小圆，周围至少5个同等大小小圆；5条直线从中心向外均匀分布（左右上下象限均覆盖）
            if len(outer_dots) < 5:
                return False
            directions = set()
            for d in outer_dots:
                dx, dy = center(d)[0] - center(center_dot)[0], center(d)[1] - center(center_dot)[1]
                if abs(dx) > abs(dy):
                    directions.add('right' if dx > 0 else 'left')
                else:
                    directions.add('down' if dy > 0 else 'up')
            return len(directions) >= 4

        def has_magnifier_and_polyline(circle):
            parts = icon_parts(circle)
            rings = [s for s in parts if get_shape_prst(s) in ('ellipse', 'arc') and 0.22 <= cm(s.width) <= 0.55 and 0.22 <= cm(s.height) <= 0.55]
            diagonal_handles = [s for s in parts if is_line(s) and abs(cm(s.width)) >= 0.12 and abs(cm(s.height)) >= 0.12]
            chart_lines = [s for s in parts if is_line(s) and line_len(s) >= 0.12]
            # 放大镜：圆/弧形镜框 + 斜向手柄；折线：至少2段可编辑线段
            return bool(rings) and bool(diagonal_handles) and len(chart_lines) >= 3

        def has_clipboard_and_cross(circle):
            parts = icon_parts(circle)
            boards = [s for s in parts if get_shape_prst(s) in ('rect', 'roundRect')
                      and 0.30 <= cm(s.width) <= 0.75 and 0.40 <= cm(s.height) <= 0.85]
            lines = [s for s in parts if is_line(s)]
            verticals = [s for s in lines if abs(cm(s.width)) <= 0.08 and cm(s.height) >= 0.12]
            horizontals = [s for s in lines if abs(cm(s.height)) <= 0.08 and cm(s.width) >= 0.12]
            # 直板夹：可编辑矩形/圆角矩形主体；医疗十字：中心一横一竖可编辑线条交叉
            return bool(boards) and bool(verticals) and bool(horizontals)

        expected = [
            (1.5, has_wireless_and_wave),
            (9.3, has_five_node_network),
            (16.5, has_magnifier_and_polyline),
            (23.8, has_clipboard_and_cross),
        ]
        for left, icon_ok in expected:
            circle = None
            for s in shapes:
                if (get_shape_prst(s) == 'ellipse'
                        and left - 0.2 <= cm(s.left) <= left + 0.3
                        and 1.0 <= cm(s.top) <= 1.6
                        and 0.7 <= cm(s.width) <= 1.2
                        and 0.7 <= cm(s.height) <= 1.2
                        and color_is_white_or_light(get_fill_color(s))):
                    circle = s
                    break
            if not circle or not icon_ok(circle):
                return False
        return True

    score, desc = (5, "+5: 流程色块图标") if check_s2_flow_icons() else (0, "+5: 流程色块图标")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: "传统机器学习"模块
    def check_s2_ml_module():
        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def contains(outer, inner, margin=0.12):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def color_is_dark_green(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return g > r and g > b and 70 <= g <= 170 and r <= 90 and b <= 110
            except:
                return False

        def color_is_light_green_fill(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return g >= r and g >= b and r >= 160 and g >= 190 and b >= 140
            except:
                return False

        def text_color_ok(color, white=False):
            if not color:
                return False
            try:
                r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
                return r >= 240 and g >= 240 and b >= 240 if white else r <= 40 and g <= 40 and b <= 40
            except:
                return False

        def text_style_ok(shape, text, size_lo, size_hi, white=False):
            if text not in get_text(shape):
                return False
            try:
                has_text = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_text = True
                        f = run.font
                        name = f.name or ""
                        size_pt = f.size / 12700 if f.size else None
                        color = str(f.color.rgb) if f.color.type else None
                        if not (name in ("Microsoft YaHei", "微软雅黑", "SimHei", "黑体")
                                and size_pt is not None and size_lo <= size_pt <= size_hi
                                and text_color_ok(color, white=white)):
                            return False
                return has_text
            except:
                return False

        flow_texts = find_shapes_with_text(slide2_shapes, "数据获取与治理")
        flow_bottom = min((shape_bottom(s) for s in flow_texts), default=2.6)

        panels = [s for s in slide2_shapes
                  if get_shape_prst(s) == 'roundRect'
                  and 2.5 <= cm(s.left) <= 8.5 and 3.0 <= cm(s.top) <= 7.5
                  and cm(s.top) >= flow_bottom
                  and 4.5 <= cm(s.width) <= 5.5 and 3.5 <= cm(s.height) <= 4.0
                  and color_is_dark_green(get_line_color(s))
                  and color_is_light_green_fill(get_fill_color(s))
                  and 0.75 <= (get_line_width_pt(s) or 0) <= 1.25]
        for panel in panels:
            title_bars = [s for s in slide2_shapes
                          if s != panel
                          and get_shape_prst(s) == 'roundRect'
                          and contains(panel, s)
                          and abs(cm(s.left) - cm(panel.left)) <= 0.2
                          and abs(cm(s.width) - cm(panel.width)) <= 0.35
                          and cm(s.top) <= cm(panel.top) + 0.35
                          and 0.7 <= cm(s.height) <= 1.0
                          and color_is_dark_green(get_fill_color(s))]
            if not title_bars:
                continue
            title_ok = False
            for bar in title_bars:
                for ts in find_shapes_with_text(slide2_shapes, "传统机器学习"):
                    if contains(bar, ts, margin=0.15) and text_style_ok(ts, "传统机器学习", 12, 14, white=True):
                        title_ok = True
                        break
                if title_ok:
                    break
            if not title_ok:
                continue

            item_ok = True
            for item_text in ("随机森林", "支持向量机"):
                item_boxes = [s for s in slide2_shapes
                              if get_shape_prst(s) == 'roundRect'
                              and contains(panel, s)
                              and cm(s.top) >= cm(panel.top) + 1.0
                              and 0.45 <= cm(s.height) <= 0.9
                              and 2.8 <= cm(s.width) <= 4.8
                              and color_is_light_green_fill(get_fill_color(s))]
                matched = False
                for box in item_boxes:
                    for ts in find_shapes_with_text(slide2_shapes, item_text):
                        if contains(box, ts, margin=0.15) and text_style_ok(ts, item_text, 11, 13, white=False):
                            matched = True
                            break
                    if matched:
                        break
                if not matched:
                    item_ok = False
                    break
            if item_ok:
                return True
        return False

    score, desc = (3, "+3: 传统机器学习模块") if check_s2_ml_module() else (0, "+3: 传统机器学习模块")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: "时序深度学习"模块
    def check_s2_dl_module():
        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def contains(outer, inner, margin=0.12):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def color_is_deep_blue(c):
            """深蓝色：b>r, b>g, 整体偏暗"""
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return b > r and b > g and (r + g + b) < 450
            except:
                return False

        def color_is_light_blue_fill(c):
            """浅蓝色填充：b>=r, b>=g, 整体偏亮"""
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return b >= r and b >= g and (r + g + b) > 500
            except:
                return False

        def text_color_ok(color, white=False):
            if not color:
                return False
            try:
                r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
                return r >= 240 and g >= 240 and b >= 240 if white else r <= 40 and g <= 40 and b <= 40
            except:
                return False

        def text_style_ok(shape, text, size_lo, size_hi, white=False):
            if text not in get_text(shape):
                return False
            try:
                has_text = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_text = True
                        f = run.font
                        name = f.name or ""
                        size_pt = f.size / 12700 if f.size else None
                        color = str(f.color.rgb) if f.color.type else None
                        if not (name in ("Microsoft YaHei", "微软雅黑", "SimHei", "黑体")
                                and size_pt is not None and size_lo <= size_pt <= size_hi
                                and text_color_ok(color, white=white)):
                            return False
                return has_text
            except:
                return False

        # 支持向量机文本底部作为模块上界参考
        svm_texts = find_shapes_with_text(slide2_shapes, "支持向量机")
        svm_bottom = min((shape_bottom(s) for s in svm_texts), default=6.5)

        panels = [s for s in slide2_shapes
                  if get_shape_prst(s) == 'roundRect'
                  and 2.5 <= cm(s.left) <= 8.5 and 6.9 <= cm(s.top) <= 11.5
                  and cm(s.top) >= svm_bottom
                  and 4.5 <= cm(s.width) <= 5.5 and 3.5 <= cm(s.height) <= 4.0
                  and color_is_deep_blue(get_line_color(s))
                  and color_is_light_blue_fill(get_fill_color(s))
                  and 0.75 <= (get_line_width_pt(s) or 0) <= 1.25]
        for panel in panels:
            title_bars = [s for s in slide2_shapes
                          if s != panel
                          and get_shape_prst(s) == 'roundRect'
                          and contains(panel, s)
                          and abs(cm(s.left) - cm(panel.left)) <= 0.2
                          and abs(cm(s.width) - cm(panel.width)) <= 0.35
                          and cm(s.top) <= cm(panel.top) + 0.35
                          and 0.7 <= cm(s.height) <= 1.0
                          and color_is_deep_blue(get_fill_color(s))]
            if not title_bars:
                continue
            title_ok = False
            for bar in title_bars:
                for ts in find_shapes_with_text(slide2_shapes, "时序深度学习"):
                    if contains(bar, ts, margin=0.15) and text_style_ok(ts, "时序深度学习", 12, 14, white=True):
                        title_ok = True
                        break
                if title_ok:
                    break
            if not title_ok:
                continue

            item_ok = True
            for item_text in ("生理信号", "生活行为序列"):
                item_boxes = [s for s in slide2_shapes
                              if get_shape_prst(s) == 'roundRect'
                              and contains(panel, s)
                              and cm(s.top) >= cm(panel.top) + 1.0
                              and 0.45 <= cm(s.height) <= 0.9
                              and 2.8 <= cm(s.width) <= 4.8
                              and color_is_light_blue_fill(get_fill_color(s))]
                matched = False
                for box in item_boxes:
                    for ts in find_shapes_with_text(slide2_shapes, item_text):
                        if contains(box, ts, margin=0.15) and text_style_ok(ts, item_text, 11, 13, white=False):
                            matched = True
                            break
                    if matched:
                        break
                if not matched:
                    item_ok = False
                    break
            if item_ok:
                return True
        return False

    score, desc = (3, "+3: 时序深度学习模块") if check_s2_dl_module() else (0, "+3: 时序深度学习模块")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 橙色虚线说明框"利用时间依赖建模..."
    def check_s2_orange_box1():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def contains(outer, inner, margin=0.15):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def color_is_orange_line(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r > 180 and g > 80 and g < 200 and b < 100
            except:
                return False

        def is_dashed(shape):
            sp = shape._element
            ln = sp.find(f'.//{DML}ln')
            if ln is not None:
                dash = ln.find(f'{DML}prstDash')
                if dash is not None:
                    return dash.get('val') in ('dash', 'lgDash', 'sysDash', 'dashDot', 'lgDashDot', 'sysDashDot')
            return False

        def text_color_is_black(color):
            if not color:
                return False
            try:
                r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
                return r <= 40 and g <= 40 and b <= 40
            except:
                return False

        def text_style_ok(shape):
            target = "利用时间依赖建模捕捉风险状态的动态变化"
            if target not in get_text(shape):
                return False
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        f = run.font
                        name = f.name or ""
                        size_pt = f.size / 12700 if f.size else None
                        color = str(f.color.rgb) if f.color.type else None
                        if not (name in ("Microsoft YaHei", "微软雅黑", "SimHei", "黑体")
                                and size_pt is not None and 11 <= size_pt <= 13
                                and f.bold
                                and text_color_is_black(color)):
                            return False
                return has_run
            except:
                return False

        # "生活行为序列"下方
        seq_texts = find_shapes_with_text(slide2_shapes, "生活行为序列")
        seq_bottom = min((shape_bottom(s) for s in seq_texts), default=10.5)

        for box in slide2_shapes:
            l, t, w, h = cm(box.left), cm(box.top), cm(box.width), cm(box.height)
            if not (get_shape_prst(box) == 'roundRect'
                    and 2.5 <= l <= 8.5 and 10.8 <= t <= 14
                    and t >= seq_bottom
                    and 6.0 <= w <= 7.5 and 1.8 <= h <= 2.2
                    and color_is_orange_line(get_line_color(box))
                    and is_dashed(box)
                    and 0.75 <= (get_line_width_pt(box) or 0) <= 1.25):
                continue
            for ts in find_shapes_with_text(slide2_shapes, "利用时间依赖建模捕捉风险状态的动态变化"):
                if contains(box, ts) and text_style_ok(ts):
                    return True
        return False

    score, desc = (1, "+1: 橙色虚线说明框(时间依赖建模)") if check_s2_orange_box1() else (0, "+1: 橙色虚线说明框(时间依赖建模)")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 绿色竖向箭头
    def check_s2_green_arrow():
        # 橙色虚线圆角矩形说明框下方，有一个高度0.8-1.0cm的绿色竖向箭头，指向下方绿色梯形“健康数据”，且为渐变式填充
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def center_x(shape):
            return cm(shape.left) + cm(shape.width) / 2

        def get_grad_stop_colors(shape):
            colors = []
            grad_fill = shape._element.find(f'.//{DML}gradFill')
            if grad_fill is None:
                return colors
            for stop in grad_fill.findall(f'.//{DML}gs'):
                srgb = stop.find(f'.//{DML}srgbClr')
                if srgb is not None:
                    colors.append(srgb.get('val'))
            return colors

        def has_green_grad_fill(shape):
            grad_fill = shape._element.find(f'.//{DML}gradFill')
            if grad_fill is None:
                return False
            return any(color_is_green(c) for c in get_grad_stop_colors(shape))

        orange_boxes = []
        for box in slide2_shapes:
            l, t, w, h = cm(box.left), cm(box.top), cm(box.width), cm(box.height)
            if (get_shape_prst(box) == 'roundRect'
                    and 2.5 <= l <= 8.5 and 10.8 <= t <= 14
                    and 6.0 <= w <= 7.5 and 1.8 <= h <= 2.2
                    and color_is_orange(get_line_color(box))
                    and get_line_dash(box) in ('dash', 'lgDash', 'sysDash', 'dashDot', 'lgDashDot', 'sysDashDot')):
                orange_boxes.append(box)

        health_trapezoids = []
        for trap in slide2_shapes:
            if get_shape_prst(trap) != 'trapezoid' or not color_is_green(get_fill_color(trap)):
                continue
            for ts in find_shapes_with_text(slide2_shapes, "健康数据"):
                if (cm(trap.left) - 0.15 <= cm(ts.left)
                        and shape_right(ts) <= shape_right(trap) + 0.15
                        and cm(trap.top) - 0.15 <= cm(ts.top)
                        and shape_bottom(ts) <= shape_bottom(trap) + 0.15):
                    health_trapezoids.append(trap)
                    break

        for arrow in slide2_shapes:
            # 一体箭头：块状箭头预设 或 带箭头端点的 line / 直线连接符
            if not is_arrow_shape(arrow):
                continue
            l, t, w, h = cm(arrow.left), cm(arrow.top), abs(cm(arrow.width)), abs(cm(arrow.height))
            if not (0.8 <= h <= 1.0 and has_green_grad_fill(arrow)):
                continue
            arrow_bottom = t + h
            arrow_cx = l + w / 2
            for box in orange_boxes:
                if not (shape_bottom(box) <= t <= shape_bottom(box) + 1.0
                        and cm(box.left) - 0.2 <= arrow_cx <= shape_right(box) + 0.2):
                    continue
                for trap in health_trapezoids:
                    if not (arrow_bottom <= cm(trap.top) + 0.2
                            and cm(box.left) - 0.4 <= center_x(trap) <= shape_right(box) + 0.4
                            and cm(trap.left) - 0.2 <= arrow_cx <= shape_right(trap) + 0.2):
                        continue
                    return True
        return False

    score, desc = (1, "+1: 绿色竖向箭头") if check_s2_green_arrow() else (0, "+1: 绿色竖向箭头")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 绿色梯形"健康数据"
    def check_s2_health_data():
        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def center_x(shape):
            return cm(shape.left) + cm(shape.width) / 2

        def contains(outer, inner, margin=0.12):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def color_is_green_or_light_green(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return g >= r and g >= b and g > 80
            except:
                return False

        def text_style_ok(shape):
            if "健康数据" not in get_text(shape):
                return False
            try:
                has_text = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_text = True
                        f = run.font
                        name = f.name or ""
                        size_pt = f.size / 12700 if f.size else None
                        if not (name in ("Microsoft YaHei", "微软雅黑", "SimHei", "黑体")
                                and size_pt is not None and 13 <= size_pt <= 15):
                            return False
                return has_text
            except:
                return False

        def is_top_long_bottom_short_trapezoid(shape):
            # Office/WPS中“上长下短”通常保存为 trapezoid 预设；宽高比例接近等边梯形外观
            return get_shape_prst(shape) == 'trapezoid' and cm(shape.width) >= cm(shape.height) * 2.0

        trapezoids = [s for s in slide2_shapes
                      if is_top_long_bottom_short_trapezoid(s)
                      and 3 <= cm(s.left) <= 7.5 and 14 <= cm(s.top) <= 16.5
                      and 2 <= cm(s.width) <= 4.5 and 0.8 <= cm(s.height) <= 1.2
                      and color_is_green_or_light_green(get_fill_color(s))]
        for trapezoid in trapezoids:
            text_ok = False
            for ts in find_shapes_with_text(slide2_shapes, "健康数据"):
                if contains(trapezoid, ts) and text_style_ok(ts):
                    text_ok = True
                    break
            if not text_ok:
                continue

            connectors = [s for s in slide2_shapes
                          if get_shape_prst(s) == 'rect'
                          and 3 <= cm(s.left) <= 7.5 and 14 <= cm(s.top) <= 16.5
                          and 0.6 <= cm(s.width) <= 0.8 and 0.8 <= cm(s.height) <= 1.2
                          and color_is_green_or_light_green(get_fill_color(s))
                          and shape_bottom(trapezoid) - 0.05 <= cm(s.top) <= shape_bottom(trapezoid) + 0.35
                          and abs(center_x(s) - center_x(trapezoid)) <= 0.25]
            if not connectors:
                continue

            for connector in connectors:
                for model_text in find_shapes_with_text(slide2_shapes, "构建预测模型"):
                    if (shape_bottom(connector) <= cm(model_text.top) + 0.4
                            and abs(center_x(connector) - center_x(model_text)) <= 1.8):
                        return True
        return False

    score, desc = (1, "+1: 绿色梯形健康数据") if check_s2_health_data() else (0, "+1: 绿色梯形健康数据")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: 三个胶囊形(构建预测模型/模型验证与整合/可视化趋势曲线)
    def check_s2_capsules():
        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def center_x(shape):
            return cm(shape.left) + cm(shape.width) / 2

        def contains(outer, inner, margin=0.15):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def text_color_is_white(color):
            if not color:
                return False
            try:
                r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
                return r >= 240 and g >= 240 and b >= 240
            except:
                return False

        def text_style_ok(shape, text, size_lo, size_hi):
            if text not in get_text(shape):
                return False
            try:
                has_text = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_text = True
                        f = run.font
                        name = f.name or ""
                        size_pt = f.size / 12700 if f.size else None
                        color = str(f.color.rgb) if f.color.type else None
                        if not (name in ("Microsoft YaHei", "微软雅黑", "SimHei", "黑体")
                                and size_pt is not None and size_lo <= size_pt <= size_hi
                                and text_color_is_white(color)):
                            return False
                return has_text
            except:
                return False

        health_trapezoids = []
        for trap in slide2_shapes:
            if not (get_shape_prst(trap) == 'trapezoid'
                    and 3 <= cm(trap.left) <= 7.5 and 14 <= cm(trap.top) <= 16.5
                    and 2 <= cm(trap.width) <= 4.5 and 0.8 <= cm(trap.height) <= 1.2
                    and color_is_green(get_fill_color(trap))):
                continue
            for ts in find_shapes_with_text(slide2_shapes, "健康数据"):
                if contains(trap, ts):
                    health_trapezoids.append(trap)
                    break

        specs = [
            ("构建预测模型", 2.5, 8.6, 16, 19, color_is_green, 14, 16),
            ("模型验证与整合", 12, 19, 16, 19, color_is_blue, 15, 17),
            ("可视化趋势曲线", 21, 28, 16, 19, color_is_purple, 15, 17),
        ]

        matched_caps = {}
        for txt, l_lo, l_hi, t_lo, t_hi, color_ok, size_lo, size_hi in specs:
            matched = None
            for cap in slide2_shapes:
                l, t, w, h = cm(cap.left), cm(cap.top), cm(cap.width), cm(cap.height)
                if not (get_shape_prst(cap) == 'roundRect'
                        and l_lo <= l <= l_hi and t_lo <= t <= t_hi
                        and 5.0 <= w <= 5.5 and 1.8 <= h <= 2.2
                        and color_ok(get_fill_color(cap))):
                    continue
                for ts in find_shapes_with_text(slide2_shapes, txt):
                    if contains(cap, ts) and text_style_ok(ts, txt, size_lo, size_hi):
                        matched = cap
                        break
                if matched:
                    break
            if not matched:
                return False
            matched_caps[txt] = matched

        first_cap = matched_caps["构建预测模型"]
        if not health_trapezoids:
            return False
        if not any(shape_bottom(trap) <= cm(first_cap.top)
                   and abs(center_x(trap) - center_x(first_cap)) <= 1.6
                   for trap in health_trapezoids):
            return False
        return True

    score, desc = (3, "+3: 三个胶囊形模块") if check_s2_capsules() else (0, "+3: 三个胶囊形模块")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: 胶囊形左侧图标
    def check_s2_capsule_icons():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def is_white_shape(shape):
            return color_is_white_or_light(get_fill_color(shape))

        def is_line_shape(shape):
            return get_shape_prst(shape) == 'line'

        def icon_bounds_ok(shapes):
            if not shapes:
                return False
            left = min(cm(s.left) for s in shapes)
            top = min(cm(s.top) for s in shapes)
            right = max(shape_right(s) for s in shapes)
            bottom = max(shape_bottom(s) for s in shapes)
            return 0.6 <= right - left <= 0.9 and 0.7 <= bottom - top <= 1.1

        def is_left_of_capsule(shapes, capsule):
            if not shapes:
                return False
            icon_right = max(shape_right(s) for s in shapes)
            icon_top = min(cm(s.top) for s in shapes)
            icon_bottom = max(shape_bottom(s) for s in shapes)
            cap_top, cap_bottom = cm(capsule.top), shape_bottom(capsule)
            return (icon_right <= cm(capsule.left) + 0.1
                    and icon_top >= cap_top - 0.1
                    and icon_bottom <= cap_bottom + 0.1)

        def near(a, b, tol=0.12):
            return abs(a - b) <= tol

        # 找三个胶囊
        def find_capsule(l_lo, l_hi, t_lo, t_hi, color_ok):
            for cap in slide2_shapes:
                l, t, w, h = cm(cap.left), cm(cap.top), cm(cap.width), cm(cap.height)
                if (get_shape_prst(cap) == 'roundRect'
                        and l_lo <= l <= l_hi and t_lo <= t <= t_hi
                        and 5.0 <= w <= 5.5 and 1.8 <= h <= 2.2
                        and color_ok(get_fill_color(cap))):
                    return cap
            return None

        cap1 = find_capsule(2.5, 8.6, 16, 19, color_is_green)
        cap2 = find_capsule(12, 19, 16, 19, color_is_blue)
        cap3 = find_capsule(21, 28, 16, 19, color_is_purple)
        if not (cap1 and cap2 and cap3):
            return False

        # ---- 图标1：构建预测模型左侧，白色竖向堆叠三段扁平圆柱 ----
        # 扁平圆柱在OOXML中为 roundRect（高度远小于宽度），至少3个白色扁平圆角矩形竖向堆叠
        stack_parts = [s for s in slide2_shapes
                       if get_shape_prst(s) in ('rect', 'roundRect')
                       and not get_text(s)
                       and is_white_shape(s)
                       and cm(s.width) >= cm(s.height) * 1.5]
        stack_parts = [s for s in stack_parts if is_left_of_capsule([s], cap1)]
        # 按纵向坐标聚合，找到至少3个纵向紧密排列的白色扁平形状
        stack_parts.sort(key=lambda s: cm(s.top))
        stacks_ok = False
        for i in range(len(stack_parts) - 2):
            trio = stack_parts[i:i+3]
            if icon_bounds_ok(trio):
                stacks_ok = True
                break
        if not stacks_ok:
            return False

        # ---- 图标2：模型验证与整合左侧，内含对勾的白色盾牌 ----
        # 办公软件中盾牌为 shield 预设或 homePlate/pentagon 等，对勾可编辑为 line 或闪电形
        shield_shapes = [s for s in slide2_shapes
                         if get_shape_prst(s) in ('shield', 'homePlate', 'pentagon', 'teardrop', 'flowChartAlternateProcess')
                         and not get_text(s)
                         and is_white_shape(s)
                         and is_left_of_capsule([s], cap2)]
        check_marks = [s for s in slide2_shapes
                       if not get_text(s) and is_white_shape(s)
                       and get_shape_prst(s) in ('line', 'curve', 'arc', 'notchedRightArrow')
                       and is_left_of_capsule([s], cap2)]
        if not shield_shapes or not check_marks:
            return False
        combined = shield_shapes + check_marks
        if not icon_bounds_ok(combined) or not is_left_of_capsule(combined, cap2):
            return False

        # ---- 图标3：可视化趋势曲线左侧，折线+白色矩形+矩形下方等宽直线 ----
        # 矩形主体
        chart_rects = [s for s in slide2_shapes
                       if get_shape_prst(s) == 'rect'
                       and not get_text(s) and is_white_shape(s)
                       and is_left_of_capsule([s], cap3)]
        # 折线（在矩形内）
        chart_lines = [s for s in slide2_shapes
                       if is_line_shape(s) and not get_text(s)
                       and is_left_of_capsule([s], cap3)]
        if not chart_rects or len(chart_lines) < 2:
            return False
        the_rect = chart_rects[0]
        # 矩形下方等宽直线（水平线，top >= rect bottom，宽度与矩形接近）
        base_lines = [s for s in chart_lines
                      if (cm(s.top) >= shape_bottom(the_rect) - 0.05
                          and near(cm(s.width), cm(the_rect.width), tol=0.25)
                          and abs(cm(s.height)) <= 0.08)]
        # 折线位于矩形内或附近
        inner_lines = [s for s in chart_lines if s not in base_lines]
        if not base_lines or not inner_lines:
            return False
        combined3 = [the_rect] + base_lines + inner_lines
        if not icon_bounds_ok(combined3) or not is_left_of_capsule(combined3, cap3):
            return False

        return True

    score, desc = (3, "+3: 胶囊形图标") if check_s2_capsule_icons() else (0, "+3: 胶囊形图标")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 胶囊间箭头
    def check_s2_capsule_arrows():
        # “模型验证与整合”左侧到“构建预测模型”间：指向左侧的绿色渐变箭头；
        # “模型验证与整合”右侧到“可视化趋势曲线”间：指向右侧的紫色渐变箭头。
        # 两个箭头均高0.4-0.8cm，长3.0-3.5cm。
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def center_y(shape):
            return cm(shape.top) + cm(shape.height) / 2

        def find_capsule(text, l_lo, l_hi, color_ok):
            for cap in slide2_shapes:
                l, t, w, h = cm(cap.left), cm(cap.top), cm(cap.width), cm(cap.height)
                if not (get_shape_prst(cap) == 'roundRect'
                        and l_lo <= l <= l_hi and 16 <= t <= 19
                        and 5.0 <= w <= 5.5 and 1.8 <= h <= 2.2
                        and color_ok(get_fill_color(cap))):
                    continue
                for ts in find_shapes_with_text(slide2_shapes, text):
                    if (cm(cap.left) - 0.15 <= cm(ts.left)
                            and cm(ts.left) + cm(ts.width) <= shape_right(cap) + 0.15
                            and cm(cap.top) - 0.15 <= cm(ts.top)
                            and cm(ts.top) + cm(ts.height) <= cm(cap.top) + cm(cap.height) + 0.15):
                        return cap
            return None

        def grad_stop_colors(shape):
            grad_fill = shape._element.find(f'.//{DML}gradFill')
            if grad_fill is None:
                return []
            colors = []
            for stop in grad_fill.findall(f'.//{DML}gs'):
                srgb = stop.find(f'.//{DML}srgbClr')
                if srgb is not None:
                    colors.append(srgb.get('val'))
            return colors

        def has_grad_color(shape, color_ok):
            return any(color_ok(c) for c in grad_stop_colors(shape))

        def arrow_ok(arrow, direction, color_ok, left_bound, right_bound, y_ref):
            prst = get_shape_prst(arrow)
            # 块状箭头预设按方向匹配；line/直线连接符只要带箭头端点即视为一体箭头
            if direction == 'left':
                if prst in ('leftArrow', 'leftRightArrow'):
                    pass
                elif is_line_or_connector(arrow) and has_line_arrowhead(arrow):
                    pass
                else:
                    return False
            elif direction == 'right':
                if prst in ('rightArrow', 'leftRightArrow'):
                    pass
                elif is_line_or_connector(arrow) and has_line_arrowhead(arrow):
                    pass
                else:
                    return False
            l, t, w, h = cm(arrow.left), cm(arrow.top), abs(cm(arrow.width)), abs(cm(arrow.height))
            if not (3.0 <= w <= 3.5 and 0.4 <= h <= 0.8):
                return False
            if not (left_bound - 0.2 <= l and l + w <= right_bound + 0.2):
                return False
            if abs(center_y(arrow) - y_ref) > 0.8:
                return False
            # 对于 line/连接符，没有块状填充，改用线条颜色判定
            if is_line_or_connector(arrow):
                lc = get_line_color(arrow)
                if lc and color_ok(lc):
                    return True
            return has_grad_color(arrow, color_ok)

        build_cap = find_capsule("构建预测模型", 2.5, 8.6, color_is_green)
        validate_cap = find_capsule("模型验证与整合", 12, 19, color_is_blue)
        visual_cap = find_capsule("可视化趋势曲线", 21, 28, color_is_purple)
        if not (build_cap and validate_cap and visual_cap):
            return False

        y_ref = center_y(validate_cap)
        green_left_arrow = False
        purple_right_arrow = False
        for s in slide2_shapes:
            if arrow_ok(s, 'left', color_is_green, shape_right(build_cap), cm(validate_cap.left), y_ref):
                green_left_arrow = True
            if arrow_ok(s, 'right', color_is_purple, shape_right(validate_cap), cm(visual_cap.left), y_ref):
                purple_right_arrow = True
        return green_left_arrow and purple_right_arrow

    score, desc = (1, "+1: 胶囊间箭头") if check_s2_capsule_arrows() else (0, "+1: 胶囊间箭头")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 蓝色虚线大框
    def check_s2_blue_dashed_frame():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def has_no_fill(shape):
            sp_pr = shape._element.find(f'.//{DML}spPr')
            return sp_pr is not None and sp_pr.find(f'{DML}noFill') is not None

        def is_white_or_no_fill(shape):
            fc = get_fill_color(shape)
            return has_no_fill(shape) or fc is None or color_is_white_or_light(fc)

        def is_blue_dashed_line(shape):
            ln = shape._element.find(f'.//{DML}ln')
            if ln is None:
                return False
            dash = ln.find(f'{DML}prstDash')
            if dash is None or dash.get('val') not in ('dash', 'lgDash', 'sysDash', 'dashDot', 'lgDashDot', 'sysDashDot'):
                return False
            lc = get_line_color(shape)
            return color_is_blue(lc) or color_is_dark_blue(lc)

        flow_bottoms = []
        for txt in ("特征提取与建模", "模型评估与解释"):
            matched = find_shapes_with_text(slide2_shapes, txt)
            if not matched:
                return False
            flow_bottoms.append(max(shape_bottom(s) for s in matched))
        required_top = max(flow_bottoms)

        for s in slide2_shapes:
            l, t, w, h = cm(s.left), cm(s.top), cm(s.width), cm(s.height)
            if not (get_shape_prst(s) == 'roundRect'
                    and 8 <= l <= 22 and 3.0 <= t <= 17
                    and t >= required_top
                    and 11 <= w <= 14 and 12 <= h <= 14):
                continue
            if (is_blue_dashed_line(s)
                    and 0.75 <= (get_line_width_pt(s) or 0) <= 1.25
                    and is_white_or_no_fill(s)):
                return True
        return False

    score, desc = (1, "+1: 蓝色虚线大框") if check_s2_blue_dashed_frame() else (0, "+1: 蓝色虚线大框")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 蓝色圆形"风险等级"
    def check_s2_risk_circle():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def contains(outer, inner, margin=0.15):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def text_color_is_white(color):
            if not color:
                return False
            try:
                r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
                return r >= 240 and g >= 240 and b >= 240
            except:
                return False

        def text_style_ok(shape):
            if "风险等级" not in get_text(shape):
                return False
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        f = run.font
                        name = f.name or ""
                        size_pt = f.size / 12700 if f.size else None
                        color = str(f.color.rgb) if f.color.type else None
                        if not (name in ("Microsoft YaHei", "微软雅黑", "SimHei", "黑体")
                                and size_pt is not None and 17 <= size_pt <= 19
                                and f.bold
                                and text_color_is_white(color)):
                            return False
                return has_run
            except:
                return False

        blue_frames = [s for s in slide2_shapes
                       if get_shape_prst(s) == 'roundRect'
                       and 8 <= cm(s.left) <= 22 and 3.0 <= cm(s.top) <= 17
                       and 11 <= cm(s.width) <= 14 and 12 <= cm(s.height) <= 14]

        for circle in slide2_shapes:
            l, t, w, h = cm(circle.left), cm(circle.top), cm(circle.width), cm(circle.height)
            if not (get_shape_prst(circle) == 'ellipse'
                    and 13 <= l <= 17 and 5 <= t <= 9
                    and 3.0 <= w <= 3.5 and 3.0 <= h <= 3.5
                    and (color_is_blue(get_fill_color(circle)) or color_is_dark_blue(get_fill_color(circle)))):
                continue
            if not any(contains(frame, circle) for frame in blue_frames):
                continue
            for ts in find_shapes_with_text(slide2_shapes, "风险等级"):
                if contains(circle, ts) and text_style_ok(ts):
                    return True
        return False

    score, desc = (1, "+1: 蓝色圆形风险等级") if check_s2_risk_circle() else (0, "+1: 蓝色圆形风险等级")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +5: 围绕风险等级8个标签+双向箭头
    def check_s2_risk_arrows():
        # 蓝色虚线圆角矩形大框内围绕“风险等级”出现8个标签，各标签与中心圆之间用双向箭头连接；
        # 箭头长度0.8-2cm，线宽0.75-1.25磅，颜色为蓝色或灰蓝色。
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def center(shape):
            return cm(shape.left) + cm(shape.width) / 2, cm(shape.top) + cm(shape.height) / 2

        def contains(outer, inner, margin=0.15):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def color_is_blue_or_gray_blue(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                blue = b >= r and b >= g and b > 80
                gray_blue = b >= r and b >= g and abs(r - g) <= 45 and 80 <= (r + g + b) / 3 <= 210
                return blue or gray_blue
            except:
                return False

        def is_double_arrow(shape):
            # 支持 line 与直线/折线连接符，两端都必须有箭头端点
            return has_double_arrowhead(shape)

        def line_endpoints(shape):
            x1, y1 = cm(shape.left), cm(shape.top)
            x2, y2 = x1 + cm(shape.width), y1 + cm(shape.height)
            return (x1, y1), (x2, y2)

        def distance(p1, p2):
            return ((p1[0] - p2[0]) ** 2 + (p1[1] - p2[1]) ** 2) ** 0.5

        def connects_arrow_to_label_and_circle(arrow, label, circle):
            p1, p2 = line_endpoints(arrow)
            cc = center(circle)
            lc = center(label)
            return ((distance(p1, cc) <= 1.1 and distance(p2, lc) <= 1.2)
                    or (distance(p2, cc) <= 1.1 and distance(p1, lc) <= 1.2))

        frames = [s for s in slide2_shapes
                  if get_shape_prst(s) == 'roundRect'
                  and 8 <= cm(s.left) <= 22 and 3.0 <= cm(s.top) <= 17
                  and 11 <= cm(s.width) <= 14 and 12 <= cm(s.height) <= 14]
        risk_circles = [s for s in slide2_shapes
                        if get_shape_prst(s) == 'ellipse'
                        and 13 <= cm(s.left) <= 17 and 5 <= cm(s.top) <= 9
                        and 3.0 <= cm(s.width) <= 3.5 and 3.0 <= cm(s.height) <= 3.5
                        and any(contains(s, ts) for ts in find_shapes_with_text(slide2_shapes, "风险等级"))]
        if not frames or not risk_circles:
            return False

        frame = frames[0]
        circle = risk_circles[0]
        label_texts = ("连续指标", "症状量表", "异常事件识别", "生活方式因素", "AUC", "SHAP", "随访评分", "F1")
        labels = []
        for txt in label_texts:
            candidates = [s for s in find_shapes_with_text(slide2_shapes, txt) if contains(frame, s) and not contains(circle, s)]
            if not candidates:
                return False
            labels.append(candidates[0])

        arrows = []
        for s in slide2_shapes:
            if not is_double_arrow(s):
                continue
            l, t, w, h = cm(s.left), cm(s.top), cm(s.width), cm(s.height)
            length = (w * w + h * h) ** 0.5
            if not (0.8 <= length <= 2.0
                    and 0.75 <= (get_line_width_pt(s) or 0) <= 1.25
                    and color_is_blue_or_gray_blue(get_line_color(s))
                    and contains(frame, s, margin=0.3)):
                continue
            arrows.append(s)

        if len(arrows) < 8:
            return False
        for label in labels:
            if not any(connects_arrow_to_label_and_circle(arrow, label, circle) for arrow in arrows):
                return False
        return True

    score, desc = (5, "+5: 围绕风险等级的双向箭头") if check_s2_risk_arrows() else (0, "+5: 围绕风险等级的双向箭头")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +5: 8个标签(连续指标/症状量表/异常事件识别/生活方式因素/AUC/SHAP/随访评分/F1)
    def check_s2_risk_labels():
        # 四个大标签与四个小指标标签都必须满足题干形状/尺寸/位置/填充/边线要求
        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def contains(outer, inner, margin=0.15):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def color_is_light_green_fill(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return g >= r and g >= b and r >= 150 and g >= 180 and b >= 120
            except:
                return False

        def color_is_light_purple_fill(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r >= 160 and b >= 160 and g >= 120 and abs(r - b) <= 80
            except:
                return False

        def color_is_light_blue_fill(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return b >= r and b >= g and r >= 130 and g >= 160 and b >= 180
            except:
                return False

        def color_is_light_yellow_fill(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r >= 200 and g >= 180 and b >= 90 and b <= 190
            except:
                return False

        # 4 个大标签: (文本, 允许 prst, 填充色判定, l_lo, l_hi, t_lo, t_hi, w_lo, w_hi, h_lo, h_hi)
        # rubric 只对填充色有要求, 未强制边线颜色/线宽, 因此这里不再对边线做硬性检查。
        big_specs = [
            ("连续指标", ('ellipse',), color_is_light_green_fill, 10, 14, 3, 7, 2.4, 2.7, 2.4, 2.7),
            ("症状量表", ('hexagon',), color_is_light_purple_fill, 17, 20, 3, 7, 1.2, 3.0, 2.0, 2.8),
            ("异常事件识别", ('diamond', 'octagon'), color_is_light_blue_fill, 17, 20, 7.5, 11, 2.5, 3.0, 2.0, 2.8),
            ("生活方式因素", ('hexagon',), color_is_light_yellow_fill, 10, 14, 7.5, 11, 1.2, 3.0, 2.0, 2.8),
        ]
        for txt, prsts, color_ok, l_lo, l_hi, t_lo, t_hi, w_lo, w_hi, h_lo, h_hi in big_specs:
            ok = False
            for shape in slide2_shapes:
                if get_shape_prst(shape) not in prsts:
                    continue
                l, t, w, h = cm(shape.left), cm(shape.top), cm(shape.width), cm(shape.height)
                if not (l_lo <= l <= l_hi and t_lo <= t <= t_hi
                        and w_lo <= w <= w_hi and h_lo <= h <= h_hi):
                    continue
                if not color_ok(get_fill_color(shape)):
                    continue
                if any(contains(shape, ts) for ts in find_shapes_with_text(slide2_shapes, txt)):
                    ok = True
                    break
            if not ok:
                return False

        # 4 个小指标标签: 圆角矩形, 浅蓝填充 + 深蓝边线 + 常规线宽 (0.75-1.25pt),
        # 高 0.6-0.8cm 宽 1.2-2.2cm。rubric 明确要求"填充为浅蓝色，框线为深蓝色的",
        # 边线颜色是硬性条件; 线宽 rubric 未写具体值,按脚本内其它框线判定的标准范围收紧。
        for txt in ("AUC", "SHAP", "随访评分", "F1"):
            ok = False
            for box in slide2_shapes:
                if get_shape_prst(box) != 'roundRect':
                    continue
                l, t, w, h = cm(box.left), cm(box.top), cm(box.width), cm(box.height)
                if not (1.2 <= w <= 2.2 and 0.6 <= h <= 0.8):
                    continue
                if not color_is_light_blue_fill(get_fill_color(box)):
                    continue
                if not color_is_dark_blue(get_line_color(box)):
                    continue
                lw = get_line_width_pt(box)
                if lw is None or not (0.75 <= lw <= 1.25):
                    continue
                if any(contains(box, ts) for ts in find_shapes_with_text(slide2_shapes, txt)):
                    ok = True
                    break
            if not ok:
                return False
        return True

    score, desc = (5, "+5: 风险等级周围8个标签") if check_s2_risk_labels() else (0, "+5: 风险等级周围8个标签")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +5: 标签上方图标
    def check_s2_label_icons():
        # 必须为可编辑组合图标，不能用单个字符文本框代替
        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def center_x(shape):
            return cm(shape.left) + cm(shape.width) / 2

        def find_label(text):
            for ts in find_shapes_with_text(slide2_shapes, text):
                return ts
            return None

        def above(label, shape, x_margin=0.7, y_gap=0.2):
            return (cm(label.left) - x_margin <= center_x(shape) <= shape_right(label) + x_margin
                    and shape_bottom(shape) <= cm(label.top) + y_gap)

        def color_is_dark_green(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return g > r and g > b and 60 <= g <= 170 and r <= 110 and b <= 120
            except:
                return False

        def color_is_dark_purple(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r >= 80 and b >= 100 and g < r and g < b and (r + b) > 200
            except:
                return False

        def color_is_cyan(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return g >= 140 and b >= 140 and r <= 120
            except:
                return False

        def color_is_yellow_icon(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r >= 180 and g >= 140 and b <= 120
            except:
                return False

        continuous = find_label("连续指标")
        symptom = find_label("症状量表")
        event = find_label("异常事件识别")
        lifestyle = find_label("生活方式因素")
        if not (continuous and symptom and event and lifestyle):
            return False

        # “连续指标”上方：六根高度上下起伏的柱状图 + 上方跟随柱状图起伏的折线，深绿色
        bars = [s for s in slide2_shapes
                if get_shape_prst(s) == 'rect' and not get_text(s)
                and above(continuous, s)
                and 0.06 <= cm(s.width) <= 0.25 and 0.25 <= cm(s.height) <= 1.0
                and color_is_dark_green(get_fill_color(s))]
        bar_heights = sorted({round(cm(s.height), 2) for s in bars})
        green_lines = [s for s in slide2_shapes
                       if get_shape_prst(s) == 'line' and not get_text(s)
                       and above(continuous, s)
                       and color_is_dark_green(get_line_color(s))]
        if len(bars) < 6 or len(bar_heights) < 3 or len(green_lines) < 4:
            return False

        # “症状量表”上方：含三条横线的纸板夹，深紫色
        purple_shapes = [s for s in slide2_shapes
                         if not get_text(s) and above(symptom, s)
                         and get_shape_prst(s) in ('rect', 'roundRect', 'line')]
        clip_boards = [s for s in purple_shapes
                       if get_shape_prst(s) in ('rect', 'roundRect')
                       and color_is_dark_purple(get_fill_color(s) or get_line_color(s))
                       and 0.5 <= cm(s.width) <= 1.1 and 0.6 <= cm(s.height) <= 1.2]
        clip_lines = [s for s in purple_shapes
                      if get_shape_prst(s) == 'line'
                      and abs(cm(s.height)) <= 0.08 and 0.3 <= cm(s.width) <= 0.9
                      and color_is_dark_purple(get_line_color(s))]
        if not clip_boards or len(clip_lines) < 3:
            return False

        # “异常事件识别”上方：内含感叹号的圆角三角形，青色
        warning_triangles = [s for s in slide2_shapes
                             if not get_text(s) and above(event, s)
                             and get_shape_prst(s) in ('triangle', 'rtTriangle')
                             and color_is_cyan(get_fill_color(s) or get_line_color(s))]
        warning_marks = [s for s in slide2_shapes
                         if not get_text(s) and above(event, s)
                         and get_shape_prst(s) in ('line', 'rect', 'ellipse')
                         and color_is_cyan(get_fill_color(s) or get_line_color(s))]
        if not warning_triangles or len(warning_marks) < 2:
            return False

        # “生活方式因素”上方：行走的黄色小人图标，高0.8-1cm、宽1-1.2cm
        yellow_parts = [s for s in slide2_shapes
                        if not get_text(s) and above(lifestyle, s)
                        and color_is_yellow_icon(get_fill_color(s) or get_line_color(s))]
        if not yellow_parts:
            return False
        left = min(cm(s.left) for s in yellow_parts)
        top = min(cm(s.top) for s in yellow_parts)
        right = max(shape_right(s) for s in yellow_parts)
        bottom = max(shape_bottom(s) for s in yellow_parts)
        if not (1.0 <= right - left <= 1.2 and 0.8 <= bottom - top <= 1.0):
            return False
        has_head = any(get_shape_prst(s) == 'ellipse' for s in yellow_parts)
        limb_lines = [s for s in yellow_parts if get_shape_prst(s) == 'line']
        if not (has_head and len(limb_lines) >= 4):
            return False
        return True

    score, desc = (5, "+5: 标签上方图标") if check_s2_label_icons() else (0, "+5: 标签上方图标")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: 横向条形图
    def check_s2_bar_chart():
        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def contains(outer, inner, margin=0.15):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def color_is_gray_or_light(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return abs(r - g) <= 30 and abs(g - b) <= 30 and abs(r - b) <= 30 and (r + g + b) > 360
            except:
                return False

        lifestyle = None
        for ts in find_shapes_with_text(slide2_shapes, "生活方式因素"):
            lifestyle = ts
            break
        if not lifestyle:
            return False
        lifestyle_bottom = shape_bottom(lifestyle)

        chart_box = None
        for s in slide2_shapes:
            l, t, w, h = cm(s.left), cm(s.top), cm(s.width), cm(s.height)
            if not (get_shape_prst(s) == 'roundRect'
                    and 9 <= l <= 15 and 10.5 <= t <= 16
                    and t >= lifestyle_bottom
                    and 4 <= w <= 6 and 3 <= h <= 5):
                continue
            fc = get_fill_color(s)
            lc = get_line_color(s)
            if color_is_white_or_light(fc) and (lc is None or color_is_gray_or_light(lc)):
                chart_box = s
                break
        if not chart_box:
            return False

        bl, bt, bw, bh = cm(chart_box.left), cm(chart_box.top), cm(chart_box.width), cm(chart_box.height)

        title_ok = any(
            "特征重要性" in get_text(ts)
            for ts in slide2_shapes
            if cm(ts.left) >= bl - 0.2 and cm(ts.top) >= bt - 0.4 and cm(ts.top) <= bt + bh + 0.3
        )
        if not title_ok:
            return False

        bars = [s for s in slide2_shapes
                if get_shape_prst(s) == 'rect'
                and contains(chart_box, s)
                and cm(s.width) >= 0.5 and cm(s.height) <= 0.3
                and not get_text(s)]
        if len(bars) < 5:
            return False

        axis_lines = [s for s in slide2_shapes
                      if get_shape_prst(s) == 'line'
                      and contains(chart_box, s)
                      and not get_text(s)]
        tick_texts = [s for s in slide2_shapes
                      if contains(chart_box, s)
                      and get_text(s)
                      and any(c.isdigit() for c in get_text(s))]
        has_axis = len(axis_lines) >= 1 or len(tick_texts) >= 2
        if not has_axis:
            return False
        return True

    score, desc = (3, "+3: 特征重要性条形图") if check_s2_bar_chart() else (0, "+3: 特征重要性条形图")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: 折线图
    def check_s2_line_chart():
        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def contains(outer, inner, margin=0.15):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def color_is_gray_or_light(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return abs(r - g) <= 30 and abs(g - b) <= 30 and abs(r - b) <= 30 and (r + g + b) > 360
            except:
                return False

        event = None
        for ts in find_shapes_with_text(slide2_shapes, "异常事件识别"):
            event = ts
            break
        if not event:
            return False
        event_bottom = shape_bottom(event)

        chart_box = None
        for s in slide2_shapes:
            l, t, w, h = cm(s.left), cm(s.top), cm(s.width), cm(s.height)
            if not (get_shape_prst(s) == 'roundRect'
                    and 14.5 <= l <= 21.5 and 10.5 <= t <= 16
                    and t >= event_bottom
                    and 4 <= w <= 6 and 3 <= h <= 5):
                continue
            fc = get_fill_color(s)
            lc = get_line_color(s)
            if color_is_white_or_light(fc) and (lc is None or color_is_gray_or_light(lc)):
                chart_box = s
                break
        if not chart_box:
            return False

        bl, bt, bw, bh = cm(chart_box.left), cm(chart_box.top), cm(chart_box.width), cm(chart_box.height)

        # 标题或说明包含"风险随时间变化"
        title_ok = any(
            "风险随时间变化" in get_text(ts)
            for ts in slide2_shapes
            if cm(ts.left) >= bl - 0.2 and cm(ts.top) >= bt - 0.4 and cm(ts.top) <= bt + bh + 0.3
        )
        if not title_ok:
            return False

        # 坐标轴：图表内可编辑 line，或含数字的刻度文本
        axis_lines = [s for s in slide2_shapes
                      if get_shape_prst(s) == 'line' and not get_text(s) and contains(chart_box, s)]
        tick_texts = [s for s in slide2_shapes
                      if contains(chart_box, s) and get_text(s) and any(c.isdigit() for c in get_text(s))]
        if not axis_lines and len(tick_texts) < 2:
            return False

        # 3条颜色不同的趋势线必须是PPT原生折线路径（custGeom），不接受多段 line 拼接
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def is_native_polyline(shape):
            if get_text(shape) or not contains(chart_box, shape) or get_line_color(shape) is None:
                return False
            sp = shape._element
            cust_geom = sp.find('.//' + DML + 'custGeom')
            if cust_geom is None:
                return False
            # 原生折线通常包含一个 moveTo 起点和至少两个 lnTo 线段终点
            ln_tos = cust_geom.findall('.//' + DML + 'lnTo')
            if len(ln_tos) < 2:
                return False
            # 排除闭合多边形或填充图形，避免把普通装饰形状误判为趋势线
            if cust_geom.find('.//' + DML + 'close') is not None:
                return False
            return True

        trend_line_colors = {
            get_line_color(s).upper()
            for s in slide2_shapes
            if is_native_polyline(s)
        }
        if len(trend_line_colors) < 3:
            return False

        # 图例文字"低风险""中风险""高风险"，每个关键字至少有一个独立文本框在图表区域内
        for k in ("低风险", "中风险", "高风险"):
            found = any(contains(chart_box, ts) for ts in find_shapes_with_text(slide2_shapes, k))
            if not found:
                return False
        return True

    score, desc = (3, "+3: 风险随时间变化折线图") if check_s2_line_chart() else (0, "+3: 风险随时间变化折线图")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: "临床表型解析与功能关联"
    def check_s2_clinical():
        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def contains(outer, inner, margin=0.15):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def text_style_ok(shape):
            if "临床表型解析与功能关联" not in get_text(shape):
                return False
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        f = run.font
                        name = f.name or ""
                        size_pt = f.size / 12700 if f.size else None
                        color = str(f.color.rgb) if f.color.type else None
                        if not (name in ("Microsoft YaHei", "微软雅黑", "SimHei", "黑体")
                                and size_pt is not None and 14 <= size_pt <= 16
                                and f.bold
                                and color_is_blue(color)):
                            return False
                return has_run
            except:
                return False

        blue_frames = [s for s in slide2_shapes
                       if get_shape_prst(s) == 'roundRect'
                       and 8 <= cm(s.left) <= 22 and 3.0 <= cm(s.top) <= 17
                       and 11 <= cm(s.width) <= 14 and 12 <= cm(s.height) <= 14
                       and (color_is_blue(get_line_color(s)) or color_is_dark_blue(get_line_color(s)))]
        validate_caps = [s for s in slide2_shapes
                         if get_shape_prst(s) == 'roundRect'
                         and 12 <= cm(s.left) <= 19 and 16 <= cm(s.top) <= 19
                         and 5.0 <= cm(s.width) <= 5.5 and 1.8 <= cm(s.height) <= 2.2
                         and color_is_blue(get_fill_color(s))]
        if not blue_frames or not validate_caps:
            return False

        for m in find_shapes_with_text(slide2_shapes, "临床表型解析与功能关联"):
            l, t = cm(m.left), cm(m.top)
            if not (12 <= l <= 19 and 15.2 <= t <= 16.4 and text_style_ok(m)):
                continue
            if not any(contains(frame, m) and shape_bottom(frame) - 2.2 <= shape_bottom(m) <= shape_bottom(frame) + 0.2
                       for frame in blue_frames):
                continue
            if not any(shape_bottom(m) <= cm(cap.top) + 0.2 for cap in validate_caps):
                continue
            return True
        return False

    score, desc = (1, "+1: 临床表型解析与功能关联") if check_s2_clinical() else (0, "+1: 临床表型解析与功能关联")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: "训练策略"模块
    def check_s2_training():
        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def contains(outer, inner, margin=0.12):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def color_is_dark_purple(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r >= 80 and b >= 100 and g < r and g < b and (r + b) > 180
            except:
                return False

        def color_is_light_purple_fill(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r >= 180 and b >= 180 and g >= 150 and abs(r - b) <= 80 and (r + g + b) > 480
            except:
                return False

        def text_color_ok(color, white=False):
            if not color:
                return False
            try:
                r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
                return r >= 240 and g >= 240 and b >= 240 if white else r <= 40 and g <= 40 and b <= 40
            except:
                return False

        def text_style_ok(shape, text, size_lo, size_hi, white=False):
            if text not in get_text(shape):
                return False
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        f = run.font
                        name = f.name or ""
                        size_pt = f.size / 12700 if f.size else None
                        color = str(f.color.rgb) if f.color.type else None
                        if not (name in ("Microsoft YaHei", "微软雅黑", "SimHei", "黑体")
                                and size_pt is not None and size_lo <= size_pt <= size_hi
                                and text_color_ok(color, white=white)):
                            return False
                return has_run
            except:
                return False

        inter_texts = find_shapes_with_text(slide2_shapes, "干预建议生成")
        inter_bottom = min((shape_bottom(s) for s in inter_texts), default=2.6)

        panels = [s for s in slide2_shapes
                  if get_shape_prst(s) == 'roundRect'
                  and 22 <= cm(s.left) <= 28 and 3.1 <= cm(s.top) <= 7.2
                  and cm(s.top) >= inter_bottom
                  and 4.5 <= cm(s.width) <= 5.5 and 3.5 <= cm(s.height) <= 4.0
                  and color_is_dark_purple(get_line_color(s))
                  and 0.75 <= (get_line_width_pt(s) or 0) <= 1.25]
        for panel in panels:
            title_bars = [s for s in slide2_shapes
                          if s != panel
                          and get_shape_prst(s) == 'roundRect'
                          and contains(panel, s)
                          and abs(cm(s.left) - cm(panel.left)) <= 0.2
                          and abs(cm(s.width) - cm(panel.width)) <= 0.35
                          and cm(s.top) <= cm(panel.top) + 0.35
                          and 0.7 <= cm(s.height) <= 1.0
                          and color_is_dark_purple(get_fill_color(s))]
            if not title_bars:
                continue
            title_ok = False
            for bar in title_bars:
                for ts in find_shapes_with_text(slide2_shapes, "训练策略"):
                    if contains(bar, ts, margin=0.15) and text_style_ok(ts, "训练策略", 12, 14, white=True):
                        title_ok = True
                        break
                if title_ok:
                    break
            if not title_ok:
                continue

            item_ok = True
            for item_text in ("交叉验证", "外部测试"):
                item_boxes = [s for s in slide2_shapes
                              if get_shape_prst(s) == 'roundRect'
                              and contains(panel, s)
                              and cm(s.top) >= cm(panel.top) + 1.0
                              and 0.45 <= cm(s.height) <= 0.9
                              and 2.8 <= cm(s.width) <= 4.8
                              and color_is_light_purple_fill(get_fill_color(s))]
                matched = False
                for box in item_boxes:
                    for ts in find_shapes_with_text(slide2_shapes, item_text):
                        if contains(box, ts, margin=0.15) and text_style_ok(ts, item_text, 11, 13, white=False):
                            matched = True
                            break
                    if matched:
                        break
                if not matched:
                    item_ok = False
                    break
            if item_ok:
                return True
        return False

    score, desc = (3, "+3: 训练策略模块") if check_s2_training() else (0, "+3: 训练策略模块")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: "模型泛化能力"模块
    def check_s2_generalization():
        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def contains(outer, inner, margin=0.12):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def color_is_light_blue_fill(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return b >= r and b >= g and r >= 130 and g >= 160 and b >= 180
            except:
                return False

        def text_color_ok(color, white=False):
            if not color:
                return False
            try:
                r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
                return r >= 240 and g >= 240 and b >= 240 if white else r <= 40 and g <= 40 and b <= 40
            except:
                return False

        def text_style_ok(shape, text, size_lo, size_hi, white=False):
            if text not in get_text(shape):
                return False
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        f = run.font
                        name = f.name or ""
                        size_pt = f.size / 12700 if f.size else None
                        color = str(f.color.rgb) if f.color.type else None
                        if not (name in ("Microsoft YaHei", "微软雅黑", "SimHei", "黑体")
                                and size_pt is not None and size_lo <= size_pt <= size_hi
                                and text_color_ok(color, white=white)):
                            return False
                return has_run
            except:
                return False

        external_texts = find_shapes_with_text(slide2_shapes, "外部测试")
        external_bottom = min((shape_bottom(s) for s in external_texts), default=6.5)

        panels = [s for s in slide2_shapes
                  if get_shape_prst(s) == 'roundRect'
                  and 22 <= cm(s.left) <= 28 and 6.9 <= cm(s.top) <= 11.5
                  and cm(s.top) >= external_bottom
                  and 4.5 <= cm(s.width) <= 5.5 and 3.5 <= cm(s.height) <= 4.0
                  and color_is_light_blue_fill(get_line_color(s))
                  and 0.75 <= (get_line_width_pt(s) or 0) <= 1.25]
        for panel in panels:
            title_bars = [s for s in slide2_shapes
                          if s != panel
                          and get_shape_prst(s) == 'roundRect'
                          and contains(panel, s)
                          and abs(cm(s.left) - cm(panel.left)) <= 0.2
                          and abs(cm(s.width) - cm(panel.width)) <= 0.35
                          and cm(s.top) <= cm(panel.top) + 0.35
                          and 0.7 <= cm(s.height) <= 1.0
                          and color_is_light_blue_fill(get_fill_color(s))]
            if not title_bars:
                continue
            title_ok = False
            for bar in title_bars:
                for ts in find_shapes_with_text(slide2_shapes, "模型泛化能力"):
                    if contains(bar, ts, margin=0.15) and text_style_ok(ts, "模型泛化能力", 12, 14, white=True):
                        title_ok = True
                        break
                if title_ok:
                    break
            if not title_ok:
                continue

            item_ok = True
            for item_text in ("稳定可靠", "参考基准"):
                item_boxes = [s for s in slide2_shapes
                              if get_shape_prst(s) == 'roundRect'
                              and contains(panel, s)
                              and cm(s.top) >= cm(panel.top) + 1.0
                              and 0.45 <= cm(s.height) <= 0.9
                              and 2.8 <= cm(s.width) <= 4.8
                              and color_is_light_blue_fill(get_fill_color(s))]
                matched = False
                for box in item_boxes:
                    for ts in find_shapes_with_text(slide2_shapes, item_text):
                        if contains(box, ts, margin=0.15) and text_style_ok(ts, item_text, 11, 13, white=False):
                            matched = True
                            break
                    if matched:
                        break
                if not matched:
                    item_ok = False
                    break
            if item_ok:
                return True
        return False

    score, desc = (3, "+3: 模型泛化能力模块") if check_s2_generalization() else (0, "+3: 模型泛化能力模块")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 橙色虚线说明框"采用ROC曲线..."
    def check_s2_roc_box():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def contains(outer, inner, margin=0.15):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def color_is_orange_line(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r > 180 and g > 80 and g < 200 and b < 100
            except:
                return False

        def is_dashed(shape):
            ln = shape._element.find(f'.//{DML}ln')
            if ln is not None:
                dash = ln.find(f'{DML}prstDash')
                if dash is not None:
                    return dash.get('val') in ('dash', 'lgDash', 'sysDash', 'dashDot', 'lgDashDot', 'sysDashDot')
            return False

        def text_style_ok(shape):
            target = "采用ROC曲线与校准曲线综合评估模型判别与泛化性能"
            if target not in get_text(shape):
                return False
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        f = run.font
                        name = f.name or ""
                        size_pt = f.size / 12700 if f.size else None
                        if not (name in ("Microsoft YaHei", "微软雅黑", "SimHei", "黑体")
                                and size_pt is not None and 10 <= size_pt <= 12):
                            return False
                return has_run
            except:
                return False

        ref_texts = find_shapes_with_text(slide2_shapes, "参考基准")
        ref_bottom = min((shape_bottom(s) for s in ref_texts), default=10.5)

        for box in slide2_shapes:
            l, t, w, h = cm(box.left), cm(box.top), cm(box.width), cm(box.height)
            if not (get_shape_prst(box) == 'roundRect'
                    and 22 <= l <= 29 and 10.8 <= t <= 14
                    and t >= ref_bottom
                    and 6.0 <= w <= 7.5 and 1.8 <= h <= 2.2
                    and color_is_orange_line(get_line_color(box))
                    and is_dashed(box)
                    and 0.75 <= (get_line_width_pt(box) or 0) <= 1.25):
                continue
            for ts in find_shapes_with_text(slide2_shapes, "采用ROC曲线与校准曲线综合评估模型判别与泛化性能"):
                if contains(box, ts) and text_style_ok(ts):
                    return True
        return False

    score, desc = (1, "+1: ROC曲线橙色说明框") if check_s2_roc_box() else (0, "+1: ROC曲线橙色说明框")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 紫色竖向箭头
    def check_s2_purple_arrow():
        DML = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def center_x(shape):
            return cm(shape.left) + cm(shape.width) / 2

        def grad_stop_colors(shape):
            grad_fill = shape._element.find(f'.//{DML}gradFill')
            if grad_fill is None:
                return []
            colors = []
            for stop in grad_fill.findall(f'.//{DML}gs'):
                srgb = stop.find(f'.//{DML}srgbClr')
                if srgb is not None:
                    colors.append(srgb.get('val'))
            return colors

        def has_purple_grad_fill(shape):
            return any(color_is_purple(c) for c in grad_stop_colors(shape))

        orange_boxes = []
        for box in slide2_shapes:
            l, t, w, h = cm(box.left), cm(box.top), cm(box.width), cm(box.height)
            if (get_shape_prst(box) == 'roundRect'
                    and 22 <= l <= 29 and 10.8 <= t <= 14
                    and 6.0 <= w <= 7.5 and 1.8 <= h <= 2.2
                    and color_is_orange(get_line_color(box))
                    and get_line_dash(box) in ('dash', 'lgDash', 'sysDash', 'dashDot', 'lgDashDot', 'sysDashDot')):
                orange_boxes.append(box)

        followup_trapezoids = []
        for trap in slide2_shapes:
            if get_shape_prst(trap) != 'trapezoid' or not color_is_purple(get_fill_color(trap)):
                continue
            for ts in find_shapes_with_text(slide2_shapes, "随访评估"):
                if (cm(trap.left) - 0.15 <= cm(ts.left)
                        and shape_right(ts) <= shape_right(trap) + 0.15
                        and cm(trap.top) - 0.15 <= cm(ts.top)
                        and shape_bottom(ts) <= shape_bottom(trap) + 0.15):
                    followup_trapezoids.append(trap)
                    break

        for arrow in slide2_shapes:
            # 一体箭头：块状箭头预设 或 带箭头端点的 line / 直线连接符
            if not is_arrow_shape(arrow):
                continue
            l, t, w, h = cm(arrow.left), cm(arrow.top), abs(cm(arrow.width)), abs(cm(arrow.height))
            if not (0.6 <= w <= 1.0 and 0.6 <= h <= 1.0 and has_purple_grad_fill(arrow)):
                continue
            arrow_bottom = t + h
            arrow_cx = l + w / 2
            for box in orange_boxes:
                if not (shape_bottom(box) <= t <= shape_bottom(box) + 1.0
                        and cm(box.left) - 0.2 <= arrow_cx <= shape_right(box) + 0.2):
                    continue
                for trap in followup_trapezoids:
                    if (arrow_bottom <= cm(trap.top) + 0.2
                            and cm(trap.left) - 0.2 <= arrow_cx <= shape_right(trap) + 0.2):
                        return True
        return False

    score, desc = (1, "+1: 紫色竖向箭头") if check_s2_purple_arrow() else (0, "+1: 紫色竖向箭头")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +3: 紫色梯形"随访评估"
    def check_s2_followup():
        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def center_x(shape):
            return cm(shape.left) + cm(shape.width) / 2

        def contains(outer, inner, margin=0.12):
            return (cm(outer.left) - margin <= cm(inner.left)
                    and shape_right(inner) <= shape_right(outer) + margin
                    and cm(outer.top) - margin <= cm(inner.top)
                    and shape_bottom(inner) <= shape_bottom(outer) + margin)

        def color_is_purple_or_light_purple(c):
            if not c:
                return False
            try:
                r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
                return r > 60 and b > 60 and r >= g * 0.7 and b >= g * 0.8
            except:
                return False

        def text_style_ok(shape):
            if "随访评估" not in get_text(shape):
                return False
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        f = run.font
                        name = f.name or ""
                        size_pt = f.size / 12700 if f.size else None
                        color = str(f.color.rgb) if f.color.type else None
                        if not (name in ("Microsoft YaHei", "微软雅黑", "SimHei", "黑体")
                                and size_pt is not None and 13 <= size_pt <= 15):
                            return False
                        try:
                            r, g, b = int((color or '000000')[0:2], 16), int((color or '000000')[2:4], 16), int((color or '000000')[4:6], 16)
                            if not (r <= 40 and g <= 40 and b <= 40):
                                return False
                        except:
                            return False
                return has_run
            except:
                return False

        def is_top_long_bottom_short(shape):
            return get_shape_prst(shape) == 'trapezoid' and cm(shape.width) >= cm(shape.height) * 2.0

        visual_caps = [s for s in slide2_shapes
                       if get_shape_prst(s) == 'roundRect'
                       and 21 <= cm(s.left) <= 28 and 16 <= cm(s.top) <= 19
                       and 5.0 <= cm(s.width) <= 5.5 and 1.8 <= cm(s.height) <= 2.2
                       and color_is_purple(get_fill_color(s))]

        trapezoids = [s for s in slide2_shapes
                      if is_top_long_bottom_short(s)
                      and 22.5 <= cm(s.left) <= 27 and 14 <= cm(s.top) <= 16
                      and 2 <= cm(s.width) <= 4.5 and 0.8 <= cm(s.height) <= 1.2
                      and color_is_purple_or_light_purple(get_fill_color(s))]
        for trapezoid in trapezoids:
            text_ok = any(contains(trapezoid, ts) and text_style_ok(ts)
                          for ts in find_shapes_with_text(slide2_shapes, "随访评估"))
            if not text_ok:
                continue

            connectors = [s for s in slide2_shapes
                          if get_shape_prst(s) == 'rect'
                          and 22.5 <= cm(s.left) <= 27 and 14 <= cm(s.top) <= 16
                          and 0.6 <= cm(s.width) <= 0.8 and 0.8 <= cm(s.height) <= 1.2
                          and color_is_purple_or_light_purple(get_fill_color(s))
                          and shape_bottom(trapezoid) - 0.05 <= cm(s.top) <= shape_bottom(trapezoid) + 0.35
                          and abs(center_x(s) - center_x(trapezoid)) <= 0.25]
            if not connectors:
                continue

            for connector in connectors:
                for cap in visual_caps:
                    if (shape_bottom(connector) <= cm(cap.top) + 0.3
                            and abs(center_x(connector) - center_x(cap)) <= 1.8):
                        return True
        return False

    score, desc = (3, "+3: 紫色梯形随访评估") if check_s2_followup() else (0, "+3: 紫色梯形随访评估")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # +1: 底部图题
    def check_s2_title():
        target = "图1 多源生理信号慢病风险预测与个体化干预研究框架"

        def shape_right(shape):
            return cm(shape.left) + cm(shape.width)

        def shape_bottom(shape):
            return cm(shape.top) + cm(shape.height)

        def center_x(shape):
            return cm(shape.left) + cm(shape.width) / 2

        def text_color_is_black(color):
            if not color:
                return False
            try:
                r, g, b = int(color[0:2], 16), int(color[2:4], 16), int(color[4:6], 16)
                return r <= 40 and g <= 40 and b <= 40
            except:
                return False

        def title_style_ok(shape):
            if target not in get_text(shape):
                return False
            try:
                has_run = False
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if not run.text.strip():
                            continue
                        has_run = True
                        f = run.font
                        name = f.name or ""
                        size_pt = f.size / 12700 if f.size else None
                        color = str(f.color.rgb) if f.color.type else None
                        if not (name in ("Microsoft YaHei", "微软雅黑", "SimHei", "黑体", "SimSun", "宋体")
                                and size_pt is not None and 13 <= size_pt <= 15
                                and f.bold
                                and text_color_is_black(color)):
                            return False
                return has_run
            except:
                return False

        validate_caps = [s for s in slide2_shapes
                         if get_shape_prst(s) == 'roundRect'
                         and 12 <= cm(s.left) <= 19 and 16 <= cm(s.top) <= 19
                         and 5.0 <= cm(s.width) <= 5.5 and 1.8 <= cm(s.height) <= 2.2
                         and color_is_blue(get_fill_color(s))]
        if not validate_caps:
            return False

        slide_center_x = cm(slides[1].slide_width) / 2
        for m in find_shapes_with_text(slide2_shapes, target):
            t = cm(m.top)
            if not (18.5 <= t <= 20 and title_style_ok(m)):
                continue
            if abs(center_x(m) - slide_center_x) > 0.8:
                continue
            if not any(shape_bottom(cap) <= t + 0.2 for cap in validate_caps):
                continue
            return True
        return False

    score, desc = (1, "+1: 第2页底部图题") if check_s2_title() else (0, "+1: 第2页底部图题")
    dim2_results.append((score, desc, score > 0))
    total_score += score

    # 汇总并返回结构化结果
    result = _build_success_result(dim1_results, dim2_results, total_score)
    result["file_name"] = file_name
    return result


if __name__ == '__main__':
    # 本地调试入口：允许通过命令行传入脚本所在目录（默认取脚本自身所在目录）
    dir_arg = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(dir_arg), ensure_ascii=False, indent=2))
