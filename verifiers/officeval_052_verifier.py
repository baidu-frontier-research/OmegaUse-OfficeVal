#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""PPTX 自动评估脚本 — 数字化校园阅读项目实施三维协同模型

统一接口约定（见"脚本接口差异与统一建议.md" §2）：
- 对外只暴露 `evaluate(dir_path: str) -> dict`
- 参数为脚本所在目录的路径；脚本自己在该目录里定位 .pptx
- 返回结构化字典（含维度一通过与否、维度二逐项得分、总分）
- 不 print 主结果、不改 sys.stdout、不 sys.exit、不硬编码路径
"""
import sys, os, re, json
from decimal import Decimal, InvalidOperation
from lxml import etree

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

SCRIPT_ID = "052"
PREFERRED_DOC_NAME = "数字化校园阅读项目实施三维协同模型_橙色可编辑版.pptx"

NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

# ─────────────────── 单位换算 ───────────────────
EMU_PER_CM = 360000
EMU_PER_PT = 12700

def cm(emu): return round(emu / EMU_PER_CM, 3)
def pt(emu): return round(emu / EMU_PER_PT, 2)


def normalize_integral_coordinates(prs):
    """规范化整数值小数，并返回不符合整数坐标类型的原始值。"""
    invalid_values = []
    node_attributes = (
        ('.//a:off | .//a:chOff | .//a:pt', ('x', 'y')),
        ('.//a:ext | .//a:chExt | .//p:sldSz', ('cx', 'cy')),
        ('.//a:path', ('w', 'h')),
    )
    for part in prs.part.package.iter_parts():
        root = getattr(part, '_element', None)
        if root is None:
            continue
        for xpath, attr_names in node_attributes:
            for node in root.xpath(xpath):
                for attr_name in attr_names:
                    raw_value = node.get(attr_name)
                    if raw_value is None:
                        continue
                    try:
                        value = Decimal(raw_value)
                    except InvalidOperation:
                        invalid_values.append(raw_value)
                        continue
                    if not value.is_finite() or value != value.to_integral_value():
                        invalid_values.append(raw_value)
                        continue
                    normalized = str(int(value))
                    if normalized != raw_value:
                        node.set(attr_name, normalized)
    return invalid_values


# ─────────────────── 几何 ───────────────────
def in_range(v, lo, hi): return lo <= v <= hi
def in_box(left_cm, top_cm, lo_l, hi_l, lo_t, hi_t):
    return in_range(left_cm, lo_l, hi_l) and in_range(top_cm, lo_t, hi_t)
def in_wh(w_cm, h_cm, lo_w, hi_w, lo_h, hi_h):
    return in_range(w_cm, lo_w, hi_w) and in_range(h_cm, lo_h, hi_h)

def shape_cm(s):
    return cm(s.left), cm(s.top), cm(s.width), cm(s.height)

def shape_center(s):
    return cm(s.left + s.width/2), cm(s.top + s.height/2)

# ─────────────────── 颜色工具 ───────────────────
def rgb_from_hex(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2],16) for i in (0,2,4))

def is_white_or_near(rgb):
    # 非主题色（白色）：只要视觉上接近白色即可通过，不要求恰为纯白。
    if rgb is None: return False
    return all(c >= 200 for c in rgb)

def is_orange(rgb):
    if rgb is None: return False
    r,g,b = rgb
    return r >= 180 and g >= 80 and b <= 80 and r > g > b

def is_deep_orange(rgb):
    if rgb is None: return False
    r,g,b = rgb
    return r >= 180 and 60 <= g <= 140 and b <= 50

def is_light_orange(rgb):
    if rgb is None: return False
    r,g,b = rgb
    return r >= 200 and g >= 120 and b <= 100 and r > g

def is_dark(rgb):
    # 非主题色（深色/黑色）：只要视觉上属于深色范围即可，不限制到某个具体主题深色。
    if rgb is None: return False
    r,g,b = rgb
    # 通道近似（避免强饱和彩色被判为深色），且整体亮度较低
    return max(rgb) <= 128 and (max(rgb) - min(rgb)) <= 60

def is_gray(rgb):
    # 非主题色（灰色）：只要是灰色调即可通过，不再限制亮度上限。
    # 三通道近似相等（视觉中性），且不是纯白（>=248 视为白）。
    if rgb is None: return False
    r,g,b = rgb
    if max(rgb) >= 248 and (max(rgb) - min(rgb)) <= 5:
        return False  # 视为白色，非灰色
    return abs(r-g) <= 30 and abs(g-b) <= 30 and abs(r-b) <= 30

def rgb_close(a, b, tol=40):
    if a is None or b is None: return False
    return all(abs(x-y)<=tol for x,y in zip(a,b))

def _xml_first_rgb(elem):
    """从 XML 元素中提取第一个 srgbClr"""
    found = elem.findall('.//{%s}srgbClr' % NS['a'])
    if found: return rgb_from_hex(found[0].get('val','000000'))
    return None

# ─────────────────── 颜色读取 ───────────────────
def get_fill_rgb(shape):
    """返回 (rgb_tuple, fill_type_str)；渐变返回主色调"""
    try:
        fill = shape.fill
        ft = str(fill.type)
        if 'SOLID' in ft:
            try: return tuple(fill.fore_color.rgb), 'solid'
            except: pass
        if 'GRADIENT' in ft:
            # 从 XML 取第一个 stop 颜色
            stops = shape._element.findall('.//{%s}gs' % NS['a'])
            if stops:
                rgb = _xml_first_rgb(stops[0])
                if rgb: return rgb, 'gradient'
            rgb = _xml_first_rgb(shape._element)
            return rgb, 'gradient'
        if 'BACKGROUND' in ft or ft == 'None': return None, 'none'
    except: pass
    return None, 'unknown'

def get_line_rgb(shape):
    try:
        lc = shape.line.color
        if lc.type is not None:
            return tuple(lc.rgb)
    except: pass
    # XML fallback
    ln = shape._element.find('.//{%s}ln' % NS['a'])
    if ln is not None:
        rgb = _xml_first_rgb(ln)
        return rgb
    return None

def get_line_width_pt(shape):
    """返回线宽（磅）。
    未显式设置线宽时，按 PowerPoint 默认线宽 1.0 磅处理——
    即只要"默认线宽 1.0 磅"落在调用方的期望范围内，就应视为通过。
    注意：这里不区分"有线但未指定宽度"与"没有线"，
    由调用方结合颜色（get_line_rgb）/ has_no_line 自行判断是否存在线条。
    """
    try:
        w = shape.line.width
        if w: return pt(w)
    except: pass
    ln = shape._element.find('.//{%s}ln' % NS['a'])
    if ln is not None:
        w = ln.get('w')
        if w: return pt(int(w))
    # 未显式设置线宽 → 采用 Office/WPS 默认线宽 1.0pt
    return 1.0

def has_dash(shape, dash_types=('dash','sysDash','lgDash','dashDot')):
    """检测是否有虚线"""
    ln = shape._element.find('.//{%s}ln' % NS['a'])
    if ln is None: return False
    pd = ln.find('{%s}prstDash' % NS['a'])
    if pd is not None:
        val = pd.get('val','')
        return any(d in val for d in dash_types)
    custDash = ln.find('{%s}custDash' % NS['a'])
    return custDash is not None

def has_no_line(shape):
    ln = shape._element.find('.//{%s}ln' % NS['a'])
    if ln is None: return True
    noFill = ln.find('{%s}noFill' % NS['a'])
    return noFill is not None


def _get_line_rgb_ignoring_alpha(shape):
    """读取边线颜色，忽略 alpha=0（透明）的情况——透明边线视为无边线。"""
    ln = shape._element.find('.//{%s}ln' % NS['a'])
    if ln is None:
        return None, True   # (rgb, is_transparent)
    noFill = ln.find('{%s}noFill' % NS['a'])
    if noFill is not None:
        return None, True
    sf = ln.find('.//{%s}solidFill' % NS['a'])
    if sf is not None:
        srgb = sf.find('{%s}srgbClr' % NS['a'])
        if srgb is not None:
            alpha = srgb.find('{%s}alpha' % NS['a'])
            if alpha is not None and alpha.get('val', '') == '0':
                return None, True   # alpha=0 即完全透明
            rgb = _xml_first_rgb(sf)
            return rgb, False
    rgb = get_line_rgb(shape)
    return rgb, False


def _line_is_no_or_thin_orange(shape, thin_pt_hi=1.0):
    """
    细则"无边线或极细橙色边线"：
    - 无边线：ln 不存在、含 noFill、或边线颜色 alpha=0（完全透明）
    - 极细橙色边线：边线颜色为橙色系且线宽 ≤ thin_pt_hi 磅
    """
    rgb, is_transparent = _get_line_rgb_ignoring_alpha(shape)
    if is_transparent or rgb is None:
        return True
    lw = get_line_width_pt(shape)
    if rgb and (is_orange(rgb) or is_light_orange(rgb) or is_deep_orange(rgb)):
        if lw is None or lw <= thin_pt_hi:
            return True
    return False


def _fill_is_orange_gradient(shape):
    """
    细则"浅橙色到白色的渐变或深浅橙色渐变"，也兜底接受实色橙色
    （Office/WPS 有时对颜色相近的渐变简化为实色存储）。
    - gradFill：stops 中含橙色系颜色即视为符合
    - solidFill：颜色为橙色系时兜底通过
    """
    el = shape._element
    grad = el.find('.//{%s}gradFill' % NS['a'])
    if grad is not None:
        stops = el.findall('.//{%s}gs' % NS['a'])
        colors = [_xml_first_rgb(s) for s in stops]
        colors = [c for c in colors if c]
        has_orange_stop = any(
            is_orange(c) or is_light_orange(c) or is_deep_orange(c)
            for c in colors
        )
        if has_orange_stop:
            return True
    # 实色橙色兜底
    rgb, ft = get_fill_rgb(shape)
    if ft == 'solid' and (is_orange(rgb) or is_light_orange(rgb) or is_deep_orange(rgb)):
        return True
    return False

def has_arrow(shape):
    """检测是否有箭头端点"""
    el = shape._element
    for tag in ('headEnd','tailEnd'):
        end = el.find('.//{%s}%s' % (NS['a'], tag))
        if end is not None:
            t = end.get('type','none')
            if t not in ('none',''):
                return True
    return False

def xml_rgb_values(shape):
    vals = []
    for el in shape._element.findall('.//{%s}srgbClr' % NS['a']):
        try: vals.append(rgb_from_hex(el.get('val')))
        except: pass
    return vals

def prst_geom(shape):
    et = shape._element.find('.//{%s}prstGeom' % NS['a'])
    return et.get('prst','') if et is not None else ''

def has_flip_h(shape) -> bool:
    """spPr/a:xfrm/@flipH == '1' 表示形状被水平翻转。
    用于区分 parallelogram 卡片"左侧较宽"(无翻转) 与"右侧较宽"(翻转)。"""
    xfrm = shape._element.find('.//{%s}xfrm' % NS['a'])
    if xfrm is None:
        return False
    return xfrm.get('flipH', '0') in ('1', 'true')

def has_round_corners(shape) -> bool:
    """线条转角是否为圆角。
    Office/PPT 中 <a:ln> 的 join 子元素 <a:round/>/<a:bevel/>/<a:miter/> 决定线条转角样式；
    显式 <a:miter/> 或 <a:bevel/> 为尖角/斜切，其余情况按 PPT 默认(圆角)处理。"""
    ln = shape._element.find('.//{%s}ln' % NS['a'])
    if ln is None:
        return True  # 无 ln 时按 PPT 默认(圆角)通过
    if ln.find('{%s}round' % NS['a']) is not None:
        return True
    if ln.find('{%s}miter' % NS['a']) is not None:
        return False
    if ln.find('{%s}bevel' % NS['a']) is not None:
        return False
    return True  # 未显式指定 join → PPT 默认圆角

def has_orange_tone(shape):
    return any(is_orange(c) or is_light_orange(c) or is_deep_orange(c) for c in xml_rgb_values(shape))

def has_gradient_fill(shape):
    # 优先检测真实渐变；若 PPT 用多层/双色橙色形状模拟渐变，也按视觉意图近似通过。
    if shape._element.find('.//{%s}gradFill' % NS['a']) is not None:
        return True
    oranges = [c for c in xml_rgb_values(shape) if is_orange(c) or is_light_orange(c) or is_deep_orange(c)]
    return len(oranges) >= 2

def has_real_gradient_fill(shape):
    return shape._element.find('.//{%s}gradFill' % NS['a']) is not None


def gradient_stop_rgbs(shape):
    colors = []
    for stop in shape._element.findall('.//{%s}gradFill/{%s}gsLst/{%s}gs' % (NS['a'], NS['a'], NS['a'])):
        rgb = _xml_first_rgb(stop)
        if rgb:
            colors.append(rgb)
    return colors


def has_linear_gradient_fill(shape):
    grad = shape._element.find('.//{%s}gradFill' % NS['a'])
    return grad is not None and grad.find('{%s}lin' % NS['a']) is not None


def has_deep_to_light_orange_linear_gradient(shape):
    colors = gradient_stop_rgbs(shape)
    return (
        has_linear_gradient_fill(shape) and
        len(colors) >= 2 and
        any(is_deep_orange(c) or is_orange(c) for c in colors) and
        any(is_light_orange(c) for c in colors) and
        all(is_deep_orange(c) or is_orange(c) or is_light_orange(c) for c in colors)
    )

def fill_is_orange(shape):
    rgb, _ = get_fill_rgb(shape)
    return is_orange(rgb) or is_light_orange(rgb) or is_deep_orange(rgb)

def fill_is_white(shape):
    rgb, _ = get_fill_rgb(shape)
    return is_white_or_near(rgb)

def has_no_fill(shape):
    try:
        fill = shape.fill
        ft = str(fill.type)
        if 'BACKGROUND' in ft or ft == 'None':
            return True
    except: pass
    return shape._element.find('.//{%s}noFill' % NS['a']) is not None

def line_is_orange_width(shape, lo, hi):
    rgb = get_line_rgb(shape)
    lw = get_line_width_pt(shape)
    return rgb and (is_orange(rgb) or is_light_orange(rgb) or is_deep_orange(rgb)) and lw and in_range(lw, lo, hi)

def line_is_white_width(shape, lo, hi):
    rgb = get_line_rgb(shape)
    lw = get_line_width_pt(shape)
    return rgb and is_white_or_near(rgb) and lw and in_range(lw, lo, hi)

# ─────────────────── 文本工具 ───────────────────
def shape_text(shape):
    if not shape.has_text_frame: return ''
    parts = []
    for para in shape.text_frame.paragraphs:
        t = ''.join(r.text for r in para.runs if r.text)
        if t: parts.append(t.strip())
    return '\n'.join(parts) if parts else shape.text_frame.text.strip()

def all_texts_in_slide(slide):
    result = []
    for s in slide.shapes:
        if s.has_text_frame:
            result.append((s, shape_text(s)))
    return result

def find_shape_with_text(slide, text, partial=False):
    for s in slide.shapes:
        if s.has_text_frame:
            t = shape_text(s)
            if partial and text in t: return s
            if not partial and t.strip() == text: return s
    return None

def _xml_rpr_attr(elem, attr):
    """从 XML 元素的第一个 rPr / defRPr 中读取属性。"""
    for tag in ('rPr', 'defRPr'):
        rpr = elem.find('.//{%s}%s' % (NS['a'], tag))
        if rpr is not None:
            val = rpr.get(attr)
            if val is not None:
                return val
    return None


def _xml_rpr_font(elem):
    """从 XML 中读取 latin typeface，覆盖 run-level 缺失的字体。"""
    for tag in ('rPr', 'defRPr'):
        rpr = elem.find('.//{%s}%s' % (NS['a'], tag))
        if rpr is not None:
            lat = rpr.find('{%s}latin' % NS['a'])
            if lat is not None:
                tf = lat.get('typeface', '')
                if tf:
                    return tf
    return None


def _xml_rpr_color(elem):
    """从 XML rPr/defRPr 中读取文字颜色（solidFill > srgbClr）。"""
    for tag in ('rPr', 'defRPr'):
        rpr = elem.find('.//{%s}%s' % (NS['a'], tag))
        if rpr is not None:
            sf = rpr.find('{%s}solidFill' % NS['a'])
            if sf is not None:
                rgb = _xml_first_rgb(sf)
                if rgb:
                    return rgb
    return None


def font_info(shape):
    """返回 (size_pt, bold, rgb, family, align) 或 None。
    先读 run 级属性，缺失时回退到 XML defRPr / lstStyle，
    确保在 Office / WPS 中继承的格式也能被正确读到。"""
    if not shape.has_text_frame: return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            f = run.font
            # ── 字号 ──
            try:
                sz = f.size and pt(f.size)
            except Exception:
                sz = None
            if not sz:
                raw = _xml_rpr_attr(run._r, 'sz')
                if raw:
                    try: sz = int(raw) / 100.0
                    except Exception: pass
            # ── 加粗 ──
            try:
                bold = f.bold
            except Exception:
                bold = None
            if bold is None:
                raw = _xml_rpr_attr(run._r, 'b')
                if raw is not None:
                    bold = raw not in ('0', 'false')
            # ── 颜色 ──
            try:
                color = tuple(f.color.rgb) if f.color.type else None
            except Exception:
                color = None
            if color is None:
                color = _xml_rpr_color(run._r)
            # ── 字体名 ──
            try:
                name = f.name
            except Exception:
                name = None
            if not name:
                name = _xml_rpr_font(run._r)
            # ── 对齐 ──
            try:
                align = para.alignment
            except Exception:
                align = None
            return sz, bold, color, name, align
    return None

def is_cjk_sans(name):
    if name is None: return True
    n = name.lower()
    # Microsoft YaHei 实际 font.name 可能是 "Microsoft YaHei"，黑体常见为 "SimHei"。
    return any(k in n for k in ('黑体','微软雅黑','microsoft yahei','heiti','simhei','yahei','sans','gothic','noto'))

def is_cjk_serif(name):
    if name is None: return True
    n = name.lower()
    return any(k in n for k in ('宋体','楷体','仿宋','song','serif','mingliU','simsun'))

def text_is_bold(shape):
    if not shape.has_text_frame: return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.bold: return True
    return False

def text_font_size_range(shape, lo, hi):
    if not shape.has_text_frame: return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                sz = run.font.size
                if sz and in_range(pt(sz), lo, hi): return True
            except: pass
    return False

def text_color_rgb(shape):
    if not shape.has_text_frame: return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                c = run.font.color
                if c.type: return tuple(c.rgb)
            except: pass
    return None

def text_is_centered(shape):
    if not shape.has_text_frame: return False
    for para in shape.text_frame.paragraphs:
        a = para.alignment
        if a in (PP_ALIGN.CENTER, None): return True
        ppr = para._p.find('{%s}pPr' % NS['a'])
        if ppr is not None and ppr.get('algn') == 'ctr': return True
    return False

def bounds_in_box(shape, lo_l, hi_l, lo_t, hi_t):
    l, t, w, h = shape_cm(shape)
    return in_range(l, lo_l, hi_l) and in_range(t, lo_t, hi_t) and l + w <= hi_l and t + h <= hi_t

def center_in_box(shape, lo_l, hi_l, lo_t, hi_t):
    cx, cy = shape_center(shape)
    return in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t)

def find_text_in_region(slide, text, lo_l, hi_l, lo_t, hi_t):
    for s in slide.shapes:
        if not s.has_text_frame: continue
        if shape_text(s).strip() != text: continue
        if bounds_in_box(s, lo_l, hi_l, lo_t, hi_t) or center_in_box(s, lo_l, hi_l, lo_t, hi_t):
            return s
    return None

def text_style_ok(shape, lo_size, hi_size, require_bold=True):
    fi = font_info(shape)
    if not fi: return False, None
    sz, bold, color, name, align = fi
    ok = (
        sz and in_range(sz, lo_size, hi_size) and
        (not require_bold or bool(bold)) and
        color and is_white_or_near(color) and
        is_cjk_sans(name) and
        text_is_centered(shape)
    )
    return ok, fi

def is_horizontal_line(shape, min_w, max_w, max_h=0.15):
    l, t, w, h = shape_cm(shape)
    return prst_geom(shape) == 'line' and in_range(w, min_w, max_w) and h <= max_h

def line_center(shape):
    l, t, w, h = shape_cm(shape)
    return l + w / 2, t + h / 2

def white_dot_ok(shape, lo=0.12, hi=0.22):
    l, t, w, h = shape_cm(shape)
    return is_oval_like(shape) and in_range(w, lo, hi) and in_range(h, lo, hi) and fill_is_white(shape)

def ordered_text_stack(shapes, max_center_delta=0.4):
    if len(shapes) < 2: return True
    centers = [shape_center(s) for s in shapes]
    xs = [c[0] for c in centers]
    ys = [c[1] for c in centers]
    return max(xs) - min(xs) <= max_center_delta and all(ys[i] < ys[i+1] for i in range(len(ys)-1))

# ─────────────────── 形状分类 ───────────────────
def is_picture(shape):
    try: return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except: return False

def is_group(shape):
    try: return shape.shape_type == MSO_SHAPE_TYPE.GROUP
    except: return False

def is_thin_line(shape):
    """窄长线形对象 (近似连接线)"""
    w, h = cm(shape.width), cm(shape.height)
    return (w <= 0.3 and h >= 1.0) or (h <= 0.3 and w >= 1.0)

def is_oval_like(shape):
    et = shape._element.find('.//{%s}prstGeom' % NS['a'])
    if et is not None:
        prst = et.get('prst','')
        return prst in ('ellipse',)
    return False

def is_rounded_rect(shape):
    et = shape._element.find('.//{%s}prstGeom' % NS['a'])
    if et is not None:
        prst = et.get('prst','')
        return prst in ('roundRect',)
    return False

# ─────────────────── shapes 查找 ───────────────────
def shapes_in_box(slide, lo_l, hi_l, lo_t, hi_t):
    result = []
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        if in_range(l, lo_l, hi_l) and in_range(t, lo_t, hi_t):
            result.append(s)
    return result

def shapes_center_in_box(slide, lo_l, hi_l, lo_t, hi_t):
    result = []
    for s in slide.shapes:
        cx, cy = shape_center(s)
        if in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t):
            result.append(s)
    return result


# ─────────────────── 维度1 ───────────────────
def check_dim1(prs, filepath):
    """返回 (passed, reasons)

    按用户指示已删除以下 rubric 项对应的检查(仅保留扩展名/幻灯片数量)：
      - 所有可见元素均为可编辑对象(图片/组合/形状数量/文本形状数量等 → 已删除)
      - 页面比例横向16:9(比例/横向检查 → 已删除)
      - 无 >50% 空白 / 无对象超出页面边界(覆盖面积/越界检查 → 已删除)
      - "核心文本存在性"锚点(实施枢纽/目标层/成效层 → 已删除)
    对应地, 维度2 中 D2-52 "整体排版完整可读"评分项也已按后续指示一并删除。
    """
    reasons = []
    ext = os.path.splitext(filepath)[1].lower()
    if ext != '.pptx':
        reasons.append(f"文件扩展名 {ext} 非 .pptx")

    if len(prs.slides) != 1:
        reasons.append(f"幻灯片数量 {len(prs.slides)} ≠ 1")

    return len(reasons) == 0, reasons


# ─────────────────── 维度2 规则定义 ───────────────────
# 每条规则: (id, score, title, checker(slide, prs))
RULES = []

def rule(rid, score, title):
    def decorator(fn):
        RULES.append((rid, score, title, fn))
        return fn
    return decorator


def is_white_or_very_light_gray_white(rgb):
    """D2-01 专用：白色或极浅灰白色，避免把浅彩色当作页面背景。"""
    if rgb is None:
        return False
    return min(rgb) >= 242 and max(rgb) - min(rgb) <= 12


def covers_whole_slide(shape, slide_w, slide_h, tol=0.15):
    l, t, w, h = shape_cm(shape)
    return l <= tol and t <= tol and l + w >= slide_w - tol and t + h >= slide_h - tol


def _a_fill_tag_exists(elem, tag):
    return elem.find('.//{%s}%s' % (NS['a'], tag)) is not None


def _solid_fill_is_white_or_graywhite(elem):
    solid = elem.find('.//{%s}solidFill' % NS['a'])
    if solid is None:
        return False
    rgb = _xml_first_rgb(solid)
    if rgb:
        return is_white_or_very_light_gray_white(rgb)
    scheme = solid.find('.//{%s}schemeClr' % NS['a'])
    if scheme is not None:
        return scheme.get('val', '') in ('bg1', 'lt1')
    return False


def _background_xml_issue(slide):
    sources = [('幻灯片', slide._element)]
    try:
        sources.append(('版式', slide.slide_layout._element))
    except Exception:
        pass
    try:
        sources.append(('母版', slide.slide_layout.slide_master._element))
    except Exception:
        pass

    for label, elem in sources:
        bg = elem.find('./{%s}cSld/{%s}bg' % (NS['p'], NS['p']))
        if bg is None:
            continue
        if _a_fill_tag_exists(bg, 'blipFill'):
            return f'{label}背景含图片/纹理填充'
        if _a_fill_tag_exists(bg, 'pattFill'):
            return f'{label}背景含底纹/图案填充'
        if _a_fill_tag_exists(bg, 'gradFill'):
            return f'{label}背景含渐变底纹'
        solid = bg.find('.//{%s}solidFill' % NS['a'])
        if solid is not None:
            rgb = _xml_first_rgb(solid)
            if rgb and not is_white_or_very_light_gray_white(rgb):
                return f'{label}背景为非白/非极浅灰白实色 #{"%02X%02X%02X"%rgb}'
    return None


def _whole_slide_shape_issue(shape, slide_w, slide_h):
    if not covers_whole_slide(shape, slide_w, slide_h):
        return None
    if is_picture(shape) or _a_fill_tag_exists(shape._element, 'blipFill'):
        return f'{shape.name} 为整页照片/图片背景'
    if _a_fill_tag_exists(shape._element, 'pattFill'):
        return f'{shape.name} 为整页底纹/图案背景'
    if _a_fill_tag_exists(shape._element, 'gradFill'):
        return f'{shape.name} 为整页渐变底纹'
    rgb, fill_type = get_fill_rgb(shape)
    if fill_type in ('solid', 'gradient') and rgb and not is_white_or_very_light_gray_white(rgb):
        return f'{shape.name} 为整页非白/非极浅灰白色块 #{"%02X%02X%02X"%rgb}'
    return None


# ── D2-01 +1 第1页白色页面背景 ──
@rule('D2-01', +1, '第1页白色页面背景')
def _(slide, prs):
    slide_w = cm(prs.slide_width); slide_h = cm(prs.slide_height)

    issue = _background_xml_issue(slide)
    if issue:
        return False, issue

    for s in slide.shapes:
        issue = _whole_slide_shape_issue(s, slide_w, slide_h)
        if issue:
            return False, issue

    bg = slide.background
    try:
        fill = bg.fill
        ft = str(fill.type)
        if 'SOLID' in ft:
            rgb = tuple(fill.fore_color.rgb)
            if is_white_or_very_light_gray_white(rgb):
                return True, '背景覆盖整张幻灯片，填充为白色或极浅灰白色，且无底纹、照片背景或整页色块水印'
            return False, f'背景为非白/非极浅灰白实色 #{"%02X%02X%02X"%rgb}'
        if 'BACKGROUND' in ft or ft == 'None':
            return True, '默认页面背景覆盖整张幻灯片，填充为白色，且无底纹、照片背景或整页色块水印'
    except Exception:
        pass

    for s in slide.shapes:
        if not covers_whole_slide(s, slide_w, slide_h):
            continue
        rgb, fill_type = get_fill_rgb(s)
        if fill_type == 'solid' and is_white_or_very_light_gray_white(rgb):
            return True, f'{s.name} 覆盖整张幻灯片，填充为白色或极浅灰白色，且无底纹、照片背景或整页色块水印'
    return False, '未检测到覆盖整张幻灯片的白色或极浅灰白色页面背景'


def top_arrow_body_ok(s):
    l, t, w, h = shape_cm(s)
    return (
        prst_geom(s) == 'upArrow' and
        in_box(l, t, 12, 21.8, 0.1, 6.8) and
        in_wh(w, h, 5.5, 8.5, 5.5, 6.8) and
        l + w <= 21.8 and t + h <= 6.8 and
        has_deep_to_light_orange_linear_gradient(s) and
        line_is_orange_width(s, 1.0, 1.8) and
        not has_dash(s)
    )

def find_top_arrow_body(slide):
    return next((s for s in slide.shapes if top_arrow_body_ok(s)), None)

def bottom_arrow_body_ok(s):
    """
    D2-07：顶部水平、底部尖角的立体下箭头（downArrow），
    位于 l:12–21.8, t:11.5–18，宽5.5–8.5cm，高5.5–6.8cm，
    深橙色到浅橙色线性渐变填充，橙色单实线边框1–1.8磅。
    """
    l, t, w, h = shape_cm(s)
    return (
        prst_geom(s) == 'downArrow' and
        in_box(l, t, 12, 21.8, 11.5, 18) and
        in_wh(w, h, 5.5, 8.5, 5.5, 6.8) and
        l + w <= 21.8 and t + h <= 18 and
        has_deep_to_light_orange_linear_gradient(s) and
        line_is_orange_width(s, 1.0, 1.8) and
        not has_dash(s)
    )


def find_bottom_arrow_body(slide):
    return next((s for s in slide.shapes if bottom_arrow_body_ok(s)), None)


def title_in_arrow_ok(slide, text, lo_l, hi_l, lo_t, hi_t):
    s = find_text_in_region(slide, text, lo_l, hi_l, lo_t, hi_t)
    if not s: return False, f'未找到位于箭头内部的"{text}"文本'
    ok, fi = text_style_ok(s, 21, 23)
    if ok:
        sz, bold, color, name, align = fi
        return True, f'{s.name} 字号{sz}pt 加粗 白色 居中'
    return False, f'"{text}"字体条件未满足: {fi}'


def top_arrow_title_ok(slide, arrow):
    al, at, aw, ah = shape_cm(arrow)
    hi_t_upper = at + ah * 0.55
    s = find_text_in_region(slide, '目标层', al, al + aw, at, hi_t_upper)
    if not s:
        return False, '未找到位于顶部中间上箭头主体内部上半部的"目标层"文本'
    fi = font_info(s)
    if not fi:
        return False, '"目标层"无法读取字体信息'
    sz, bold, color, name, align = fi
    if not (sz and in_range(sz, 21, 23)):
        return False, f'"目标层"字号{sz}pt 不在21–23磅范围'
    if not bold:
        return False, '"目标层"未加粗'
    if not (color and is_white_or_near(color)):
        return False, f'"目标层"颜色{color}非白色'
    if not is_cjk_sans(name):
        return False, f'"目标层"字体"{name}"非黑体/微软雅黑或相近无衬线字体'
    if not text_is_centered(s):
        return False, '"目标层"未水平居中'
    return True, f'{s.name} 位于上箭头上半部 字号{sz}pt 加粗 白色 居中 字体{name}'


def bottom_arrow_title_ok(slide, arrow):
    al, at, aw, ah = shape_cm(arrow)
    # 下箭头主体内部上半部：顶部到主体高度约 55% 处
    hi_t_upper = at + ah * 0.55
    s = find_text_in_region(slide, '成效层', al, al + aw, at, hi_t_upper)
    if not s:
        return False, '未找到位于底部中间下箭头主体内部上半部的"成效层"文本'
    fi = font_info(s)
    if not fi:
        return False, '"成效层"无法读取字体信息'
    sz, bold, color, name, align = fi
    if not (sz and in_range(sz, 21, 23)):
        return False, f'"成效层"字号{sz}pt 不在21–23磅范围'
    if not bold:
        return False, '"成效层"未加粗'
    if not (color and is_white_or_near(color)):
        return False, f'"成效层"颜色{color}非白色'
    if not is_cjk_sans(name):
        return False, f'"成效层"字体"{name}"非黑体/微软雅黑或相近无衬线字体'
    if not text_is_centered(s):
        return False, '"成效层"未水平居中'
    return True, f'{s.name} 位于下箭头上半部 字号{sz}pt 加粗 白色 居中 字体{name}'

def separator_and_dot_ok(slide, lo_l, hi_l, lo_t, hi_t):
    lines = []
    dots = []
    for s in slide.shapes:
        if not (bounds_in_box(s, lo_l, hi_l, lo_t, hi_t) or center_in_box(s, lo_l, hi_l, lo_t, hi_t)): continue
        l, t, w, h = shape_cm(s)
        if is_horizontal_line(s, 4, 6, 0.12) and line_is_white_width(s, 0.75, 1.25):
            lines.append(s)
        if white_dot_ok(s, 0.12, 0.22):
            dots.append(s)
    for line in lines:
        lx, ly = line_center(line)
        for dot in dots:
            dx, dy = shape_center(dot)
            if abs(dx - lx) <= 0.2 and 0.0 <= dy - ly <= 0.35:
                return True, '找到中心对齐的白色分隔线和白色小圆点'
    return False, f'分隔线{len(lines)}条 小圆点{len(dots)}个，未满足中心对齐/上下关系'


def _find_title_shape(slide, text, lo_l, hi_l, lo_t, hi_t):
    """在给定区域内找到精确文本匹配的形状。"""
    return find_text_in_region(slide, text, lo_l, hi_l, lo_t, hi_t)


def _arrow_sep_dot_ok(slide, arrow, title_text):
    """
    通用：顶部/底部箭头内分隔线与白色小圆点检查。
    1. 找到 title_text 文本形状，取其底边作为参考 Y。
    2. 分隔线须在参考 Y 下方 0.3–0.6 cm 处（中心 Y）。
    3. 分隔线：白色水平单实线，长 4–6 cm，线宽 0.75–1.25 磅，无虚线。
    4. 分隔线下方中间须有一个白色小圆点，直径 0.12–0.22 cm。
       "下方中间"：圆点中心 X 与分隔线中心 X 偏差 ≤ 0.2 cm，
                  圆点中心 Y 在分隔线中心 Y 到 +0.35 cm 之间。
    """
    al, at, aw, ah = shape_cm(arrow)
    title_s = _find_title_shape(slide, title_text, al, al + aw, at, at + ah)
    if not title_s:
        return False, f'未找到"{title_text}"文本，无法确定分隔线参考位置'

    _, title_top, _, title_h = shape_cm(title_s)
    title_bottom = title_top + title_h

    lines = []
    dots  = []
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        cx, cy = shape_center(s)
        if not (in_range(cx, al, al + aw) and in_range(cy, at, at + ah)):
            continue
        # ── 分隔线：白色水平单实线，长 4–6 cm，线宽 0.75–1.25 磅 ──
        if is_horizontal_line(s, 4, 6, 0.12) and line_is_white_width(s, 0.75, 1.25) and not has_dash(s):
            if in_range(cy - title_bottom, 0.3, 0.6):
                lines.append(s)
        # ── 白色小圆点，直径 0.12–0.22 cm ──
        if white_dot_ok(s, 0.12, 0.22):
            dots.append(s)

    if not lines:
        return False, f'未找到位于"{title_text}"底边下方0.3–0.6cm的白色水平单实线（4–6cm，0.75–1.25磅）'

    for line in lines:
        lx, ly = line_center(line)
        for dot in dots:
            dx, dy = shape_center(dot)
            if abs(dx - lx) <= 0.2 and 0.0 <= dy - ly <= 0.35:
                l2, t2, w2, h2 = shape_cm(line)
                d2l, d2t, d2w, d2h = shape_cm(dot)
                return True, (
                    f'分隔线{line.name} 长{w2:.2f}cm 位于"{title_text}"下方{ly-title_bottom:.2f}cm；'
                    f'圆点{dot.name} 直径约{(d2w+d2h)/2:.3f}cm 在分隔线下方中间'
                )

    return False, (
        f'找到分隔线{len(lines)}条，白色小圆点{len(dots)}个，'
        f'但无圆点满足位于分隔线下方中间（X偏差≤0.2cm，Y偏差0–0.35cm）'
    )


def top_arrow_sep_dot_ok(slide, arrow, title_text='目标层'):
    return _arrow_sep_dot_ok(slide, arrow, title_text)


def bottom_arrow_sep_dot_ok(slide, arrow, title_text='成效层'):
    return _arrow_sep_dot_ok(slide, arrow, title_text)


def four_line_texts_ok(slide, required, lo_l, hi_l, lo_t, hi_t):
    text_shapes = []
    for text in required:
        s = find_text_in_region(slide, text, lo_l, hi_l, lo_t, hi_t)
        if not s:
            return False, f'缺少箭头内部文本"{text}"'
        ok, fi = text_style_ok(s, 13, 15)
        if not ok:
            return False, f'"{text}"字体条件未满足: {fi}'
        text_shapes.append(s)
    if not ordered_text_stack(text_shapes):
        return False, '四行文本未按自上而下顺序或水平中心未对齐'

    dashes = []
    for s in slide.shapes:
        if not (bounds_in_box(s, lo_l, hi_l, lo_t, hi_t) or center_in_box(s, lo_l, hi_l, lo_t, hi_t)): continue
        if is_horizontal_line(s, 4, 6.2, 0.12) and line_is_white_width(s, 0.75, 1.25) and has_dash(s):
            dashes.append(s)
    if len(dashes) >= 3:
        return True, f'四行文本均找到，{len(dashes)}条白色虚线'
    return False, f'四行文本找到，但白色虚线仅{len(dashes)}条'


def _font_info_all_runs(shape):
    """返回 shape 内所有 run 的字体信息列表，每项 (sz, bold, rgb, name, align)。"""
    if not shape.has_text_frame:
        return []
    results = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            f = run.font
            try: sz = f.size and pt(f.size)
            except: sz = None
            if not sz:
                raw = _xml_rpr_attr(run._r, 'sz')
                if raw:
                    try: sz = int(raw) / 100.0
                    except: pass
            try: bold = f.bold
            except: bold = None
            if bold is None:
                raw = _xml_rpr_attr(run._r, 'b')
                if raw is not None:
                    bold = raw not in ('0', 'false')
            try: color = tuple(f.color.rgb) if f.color.type else None
            except: color = None
            if color is None:
                color = _xml_rpr_color(run._r)
            try: name = f.name
            except: name = None
            if not name:
                name = _xml_rpr_font(run._r)
            try: align = para.alignment
            except: align = None
            results.append((sz, bold, color, name, align))
    return results


def _is_bold_or_semibold(bold, font_name):
    if bold is True:
        return True
    if not font_name:
        return False
    n = font_name.lower()
    return any(k in n for k in ('semibold', 'semi bold', 'demibold', 'demi bold', 'medium', '中黑', '准黑'))


def _row_text_style_ok(shape):
    """
    检查一个文本框是否满足四行说明文本的样式要求：
    - 字号 13–15 磅
    - 加粗或半加粗（bold=True，或字体族名称明确为 semibold/demibold/medium/中黑/准黑）
    - 颜色为白色
    - 字体为黑体/微软雅黑或相近无衬线字体
    - 水平居中
    """
    all_runs = _font_info_all_runs(shape)
    if not all_runs:
        fi = font_info(shape)
        if not fi:
            return False, '无字体信息'
        all_runs = [fi]

    sz_ok = any(r[0] and in_range(r[0], 13, 15) for r in all_runs)
    bold_ok = any(_is_bold_or_semibold(r[1], r[3]) for r in all_runs)
    color_ok = any(r[2] and is_white_or_near(r[2]) for r in all_runs)
    font_ok = any(is_cjk_sans(r[3]) for r in all_runs)
    center_ok = text_is_centered(shape)

    if not sz_ok:
        szs = [r[0] for r in all_runs]
        return False, f'字号{szs}不在13–15磅范围'
    if not bold_ok:
        return False, '未加粗/半加粗'
    if not color_ok:
        cols = [r[2] for r in all_runs]
        return False, f'颜色{cols}非白色'
    if not font_ok:
        names = [r[3] for r in all_runs]
        return False, f'字体{names}非黑体/微软雅黑或相近无衬线字体'
    if not center_ok:
        return False, '未水平居中'
    return True, 'ok'


def top_arrow_four_lines_ok(slide, arrow,
                             required=('核心素养导向', '阅读兴趣培养', '信息获取与表达', '校园文化共建')):
    return _arrow_four_lines_ok(slide, arrow, required)


def bottom_arrow_four_lines_ok(slide, arrow,
                                required=('阅读习惯形成', '思维能力提升', '协作交流增强', '自主学习深化')):
    return _arrow_four_lines_ok(slide, arrow, required)


def _arrow_four_lines_ok(slide, arrow, required):
    """
    顶部/底部箭头内四行说明文本 + 三条白色虚线，全部相对箭头坐标定位。
    细则要点：
    1. 位于箭头内部中下部（箭头高度 45%–100% 区间）
    2. 四行文本按顺序从上到下出现
    3. 每行：白色、黑体/微软雅黑、字号 13–15 磅、加粗或半加粗、水平居中
    4. 相邻行垂直间距 0.35–0.7 cm（按文本框 top 边之差计算）
    5. 四行之间三条白色虚线，长 4–6 cm，线宽 0.75–1.25 磅
    """
    al, at, aw, ah = shape_cm(arrow)
    lo_t = at + ah * 0.45
    hi_t = at + ah
    lo_l, hi_l = al, al + aw

    # ── 1. 找四行文本并逐一验证样式 ──
    text_shapes = []
    for text in required:
        s = find_text_in_region(slide, text, lo_l, hi_l, lo_t, hi_t)
        if not s:
            return False, f'在箭头中下部未找到文本"{text}"'
        ok, reason = _row_text_style_ok(s)
        if not ok:
            return False, f'"{text}"样式不满足：{reason}'
        text_shapes.append(s)

    # ── 2. 四行自上而下排列，水平中心对齐 ──
    centers = [shape_center(s) for s in text_shapes]
    ys = [c[1] for c in centers]
    xs = [c[0] for c in centers]
    if not all(ys[i] < ys[i + 1] for i in range(len(ys) - 1)):
        return False, f'四行文本未按自上而下顺序排列，中心Y={ys}'
    if max(xs) - min(xs) > 0.5:
        return False, f'四行文本水平中心未对齐，X偏差{max(xs)-min(xs):.2f}cm'

    # ── 3. 相邻行垂直间距 0.35–0.7 cm ──
    tops = [shape_cm(s)[1] for s in text_shapes]
    for i in range(len(tops) - 1):
        gap = tops[i + 1] - tops[i]
        if not in_range(gap, 0.35, 0.7):
            return False, (
                f'"{required[i]}"与"{required[i+1]}"垂直间距{gap:.2f}cm，'
                f'不在0.35–0.7cm范围'
            )

    # ── 4. 四行之间三条白色虚线（在相邻两行文本的 Y 范围内各找一条） ──
    dashes_between = 0
    for i in range(len(text_shapes) - 1):
        y_lo = ys[i]
        y_hi = ys[i + 1]
        found = False
        for s in slide.shapes:
            cx, cy = shape_center(s)
            if not (in_range(cx, lo_l, hi_l) and in_range(cy, y_lo, y_hi)):
                continue
            if (is_horizontal_line(s, 4, 6, 0.12) and
                    line_is_white_width(s, 0.75, 1.25) and
                    has_dash(s)):
                found = True
                break
        if found:
            dashes_between += 1

    if dashes_between < 3:
        all_white_lines = [
            s for s in slide.shapes
            if (in_range(shape_center(s)[0], lo_l, hi_l) and
                in_range(shape_center(s)[1], lo_t, hi_t) and
                is_horizontal_line(s, 4, 6, 0.12) and
                line_is_white_width(s, 0.75, 1.25))
        ]
        all_dashes = [s for s in all_white_lines if has_dash(s)]
        if len(all_dashes) >= 3:
            dashes_between = 3
        else:
            return False, (
                f'四行文本间仅检测到{dashes_between}条白色虚线（每相邻行间需各一条），'
                f'区域内虚线总计{len(all_dashes)}条'
            )

    return True, (
        f'四行文本均找到且样式符合，行间距{[round(tops[i+1]-tops[i],2) for i in range(3)]}cm，'
        f'{dashes_between}条白色虚线分隔'
    )


def ellipse_ring_ok(s, lo_w, hi_w, lo_h, hi_h, lo_lw, hi_lw, dashed=False):
    l, t, w, h = shape_cm(s); cx, cy = shape_center(s)
    return (
        is_oval_like(s) and
        in_range(cx, 12, 22) and in_range(cy, 6, 12) and
        in_wh(w, h, lo_w, hi_w, lo_h, hi_h) and
        has_no_fill(s) and
        line_is_orange_width(s, lo_lw, hi_lw) and
        (has_dash(s) if dashed else not has_dash(s))
    )


def outer_ellipse_ring_ok(s):
    """
    D2-12 专用：中部"实施枢纽"圆角矩形框最外层横向椭圆环。
    逐点检查：位置范围、横向长椭圆、尺寸、无填充、浅橙色单实线、线宽。
    白色填充（rgb≥220且接近白色）也视同无填充，兼容 Office/WPS 用白色填充模拟透明的写法。
    """
    l, t, w, h = shape_cm(s)
    line_rgb = get_line_rgb(s)
    lw = get_line_width_pt(s)
    fill_rgb, fill_type = get_fill_rgb(s)
    fill_ok = has_no_fill(s) or (fill_type == 'solid' and is_white_or_near(fill_rgb))
    return (
        is_oval_like(s) and
        in_box(l, t, 9, 25, 6, 12) and
        l + w <= 25 and t + h <= 12 and
        w > h and
        in_wh(w, h, 14, 16, 5.0, 5.6) and
        fill_ok and
        line_rgb and is_light_orange(line_rgb) and
        lw and in_range(lw, 1.5, 2.2) and
        not has_dash(s)
    )


# ── D2-02 +3 顶部中间上箭头主体 ──
@rule('D2-02', +3, '顶部中间上箭头主体')
def _(slide, prs):
    s = find_top_arrow_body(slide)
    if s:
        l, t, w, h = shape_cm(s)
        stops = gradient_stop_rgbs(s)
        return True, f'{s.name} 位置({l:.1f},{t:.1f}) 尺寸{w:.1f}x{h:.1f}cm upArrow 深橙→浅橙线性渐变 橙色实线边框'
    return False, '未找到符合位置、尺寸、上箭头形状、深橙→浅橙线性渐变填充和橙色单实线边框要求的顶部上箭头主体'


# ── D2-03 +1 顶部箭头中间标题文本"目标层" ──
@rule('D2-03', +1, '顶部箭头中间标题文本"目标层"')
def _(slide, prs):
    arrow = find_top_arrow_body(slide)
    if not arrow:
        return False, '顶部上箭头主体未通过，标题不计入箭头内部文本'
    return top_arrow_title_ok(slide, arrow)


# ── D2-04 +1 顶部箭头分隔线与小圆点 ──
@rule('D2-04', +1, '顶部箭头分隔线与白色小圆点')
def _(slide, prs):
    arrow = find_top_arrow_body(slide)
    if not arrow:
        return False, '顶部上箭头主体未通过，分隔线和圆点不计分'
    return top_arrow_sep_dot_ok(slide, arrow, '目标层')


# ── D2-05 +3 顶部中间四行说明文本与虚线 ──
@rule('D2-05', +3, '顶部中间四行说明文本与虚线')
def _(slide, prs):
    arrow = find_top_arrow_body(slide)
    if not arrow:
        return False, '顶部上箭头主体未通过，四行说明文本不计分'
    return top_arrow_four_lines_ok(slide, arrow)


# ── D2-06 +1 顶部中间向下渐变箭头 ──
@rule('D2-06', +1, '顶部中间向下渐变箭头')
def _(slide, prs):
    top_arrow = find_top_arrow_body(slide)
    if not top_arrow:
        return False, '顶部上箭头主体未通过，顶部中间向下渐变箭头不计分'

    impl_frame = _find_impl_frame(slide)
    if not impl_frame:
        return False, '中间实施框未通过，无法确认向下渐变箭头位于其上方'

    top_l, top_t, top_w, top_h = shape_cm(top_arrow)
    impl_l, impl_t, impl_w, impl_h = shape_cm(impl_frame)
    top_bottom = top_t + top_h

    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        cx, cy = shape_center(s)
        if prst_geom(s) != 'downArrow':
            continue
        if not (in_box(l, t, 16, 18, 6, 8) and in_wh(w, h, 0.6, 1.6, 0.8, 1.2)):
            continue
        # 位于顶部上箭头主体下方、中心实施框上方
        if not (t >= top_bottom - 0.2 and t + h <= impl_t + 0.2):
            continue
        # 竖向箭头：高度大于宽度，且中心位于顶部上箭头和实施框的水平重叠范围内
        if not (h > w and in_range(cx, max(top_l, impl_l) - 0.5, min(top_l + top_w, impl_l + impl_w) + 0.5)):
            continue
        if not _fill_is_orange_gradient(s):
            continue
        if not _line_is_no_or_thin_orange(s, 1.0):
            continue
        return True, f'{s.name} downArrow 位置({l:.1f},{t:.1f}) 尺寸{w:.1f}x{h:.1f}cm 橙白/深浅橙渐变 无边线或极细橙边'
    return False, '未找到符合位置、朝下竖向箭头、尺寸、橙白/深浅橙渐变填充和无边线/极细橙边要求的顶部中间向下箭头'


# ── D2-07 +3 底部中间下箭头主体 ──
@rule('D2-07', +3, '底部中间下箭头主体')
def _(slide, prs):
    s = find_bottom_arrow_body(slide)
    if s:
        l, t, w, h = shape_cm(s)
        return True, f'{s.name} downArrow {w:.1f}x{h:.1f}cm 深橙→浅橙线性渐变 橙色单实线边框'

    # 兜底：rect+triangle 拼合下箭头（两部分各自满足渐变条件）
    rect = None; tri = None
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        if not in_box(l, t, 12, 21.8, 11.5, 18): continue
        if not line_is_orange_width(s, 1.0, 1.8): continue
        if not has_deep_to_light_orange_linear_gradient(s): continue
        if prst_geom(s) == 'rect' and in_range(w, 5.5, 8.5) and in_range(h, 3.0, 4.2): rect = s
        if prst_geom(s) == 'triangle' and in_range(w, 5.5, 8.5) and in_range(h, 1.5, 2.6): tri = s
    if rect and tri:
        rl, rt, rw, rh = shape_cm(rect); rcx, _ = shape_center(rect)
        tl, tt, tw, th = shape_cm(tri); tcx, _ = shape_center(tri)
        total_l = min(rl, tl); total_t = min(rt, tt)
        total_r = max(rl + rw, tl + tw); total_b = max(rt + rh, tt + th)
        if (abs(rcx - tcx) <= 0.15 and abs((rt + rh) - tt) <= 0.25 and
                in_wh(total_r - total_l, total_b - total_t, 5.5, 8.5, 5.5, 6.8) and
                in_box(total_l, total_t, 12, 21.8, 11.5, 18) and
                total_r <= 21.8 and total_b <= 18):
            return True, f'{rect.name}+{tri.name} 组合下箭头 深橙→浅橙线性渐变 橙色单实线边框'
    return False, '未找到符合位置、downArrow形状、尺寸、深橙→浅橙线性渐变填充和橙色单实线边框要求的底部下箭头主体'


# ── D2-08 +1 底部箭头内标题文本"成效层" ──
@rule('D2-08', +1, '底部箭头内标题文本"成效层"')
def _(slide, prs):
    arrow = find_bottom_arrow_body(slide)
    if not arrow:
        return False, '底部下箭头主体未通过，标题不计入箭头内部文本'
    return bottom_arrow_title_ok(slide, arrow)


# ── D2-09 +1 底部箭头内分隔线与小圆点 ──
@rule('D2-09', +1, '底部箭头内分隔线与白色小圆点')
def _(slide, prs):
    arrow = find_bottom_arrow_body(slide)
    if not arrow:
        return False, '底部下箭头主体未通过，分隔线和圆点不计分'
    return bottom_arrow_sep_dot_ok(slide, arrow, '成效层')


# ── D2-10 +3 底部中间四行说明文本与虚线 ──
@rule('D2-10', +3, '底部中间四行说明文本与虚线')
def _(slide, prs):
    arrow = find_bottom_arrow_body(slide)
    if not arrow:
        return False, '底部下箭头主体未通过，四行说明文本不计分'
    return bottom_arrow_four_lines_ok(slide, arrow)


# ── D2-11 +1 中间向下渐变箭头 ──
@rule('D2-11', +1, '中间向下渐变箭头（实施框下方）')
def _(slide, prs):
    impl_frame = _find_impl_frame(slide)
    if not impl_frame:
        return False, '中间实施框未通过，无法确认中间向下渐变箭头位于其下方'

    bottom_arrow = find_bottom_arrow_body(slide)
    if not bottom_arrow:
        return False, '底部下箭头主体未通过，无法确认中间向下渐变箭头位于其上方'

    impl_l, impl_t, impl_w, impl_h = shape_cm(impl_frame)
    bottom_l, bottom_t, bottom_w, bottom_h = shape_cm(bottom_arrow)
    impl_bottom = impl_t + impl_h

    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        cx, cy = shape_center(s)
        if prst_geom(s) != 'downArrow':
            continue
        if not (in_box(l, t, 16, 18, 10.0, 12.5) and in_wh(w, h, 1.0, 1.4, 1.5, 1.8)):
            continue
        # 位于中心实施框下方、底部中间下箭头主体上方
        if not (t >= impl_bottom - 0.2 and t + h <= bottom_t + 0.2):
            continue
        # 竖向箭头，且中心在实施框和底部箭头水平重叠区域附近
        if not (h > w and in_range(cx, max(impl_l, bottom_l) - 0.5, min(impl_l + impl_w, bottom_l + bottom_w) + 0.5)):
            continue
        if not _fill_is_orange_gradient(s):
            continue
        return True, f'{s.name} downArrow 位置({l:.1f},{t:.1f}) 尺寸{w:.1f}x{h:.1f}cm 橙白/深浅橙渐变'
    return False, '未找到符合实施框下方、底部下箭头上方、位置、竖向朝下箭头、尺寸和橙白/深浅橙渐变填充要求的中间向下箭头'


# ── D2-12 +1 最外层横向椭圆环 ──
@rule('D2-12', +1, '中部最外层横向椭圆环')
def _(slide, prs):
    for s in slide.shapes:
        if outer_ellipse_ring_ok(s):
            l, t, w, h = shape_cm(s)
            lw = get_line_width_pt(s)
            return True, f'{s.name} {w:.1f}x{h:.1f}cm 浅橙色无填充实线椭圆环 线宽{lw}pt'
    return False, '未找到符合距左9–25cm、距上6–12cm、宽14–16cm、高5–5.6cm、无填充、浅橙色单实线、线宽1.5–2.2磅要求的最外层横向椭圆环'


def second_ellipse_ring_ok(s):
    """
    D2-13 专用：中部"实施枢纽"圆角矩形框第二层横向椭圆环。
    逐点检查：位置范围、横向长椭圆、尺寸、无填充（或白色填充兜底）、浅橙色单实线、线宽。
    """
    l, t, w, h = shape_cm(s)
    line_rgb = get_line_rgb(s)
    lw = get_line_width_pt(s)
    fill_rgb, fill_type = get_fill_rgb(s)
    fill_ok = has_no_fill(s) or (fill_type == 'solid' and is_white_or_near(fill_rgb))
    return (
        is_oval_like(s) and
        in_box(l, t, 10, 24, 6, 12) and
        l + w <= 24 and t + h <= 12 and
        w > h and
        in_wh(w, h, 12.5, 13, 4.5, 5.5) and
        fill_ok and
        line_rgb and is_light_orange(line_rgb) and
        lw and in_range(lw, 1.5, 2.2) and
        not has_dash(s)
    )


# ── D2-13 +1 第二层横向椭圆环 ──
@rule('D2-13', +1, '中部第二层横向椭圆环')
def _(slide, prs):
    for s in slide.shapes:
        if second_ellipse_ring_ok(s):
            l, t, w, h = shape_cm(s)
            lw = get_line_width_pt(s)
            return True, f'{s.name} {w:.1f}x{h:.1f}cm 第二层浅橙色无填充实线椭圆环 线宽{lw}pt'
    return False, '未找到符合距左10–24cm、距上6–12cm、宽12.5–13cm、高4.5–5.5cm、无填充、浅橙色单实线、线宽1.5–2.2磅要求的第二层横向椭圆环'


def inner_ellipse_ring_ok(s):
    """
    D2-14 专用：中部"实施枢纽"圆角矩形框最内侧横向椭圆环。
    逐点检查：位于最外层内侧、位置范围、横向长椭圆、尺寸、无填充（或白色填充兜底）、
    浅橙色虚线、线宽。"虚线间距均匀"对应标准 prstDash 类型（非自定义 custDash）。
    """
    l, t, w, h = shape_cm(s)
    line_rgb = get_line_rgb(s)
    lw = get_line_width_pt(s)
    fill_rgb, fill_type = get_fill_rgb(s)
    fill_ok = has_no_fill(s) or (fill_type == 'solid' and is_white_or_near(fill_rgb))
    # 虚线且间距均匀：prstDash（标准预设）而非 custDash（自定义）
    ln = s._element.find('.//{%s}ln' % NS['a'])
    dash_ok = False
    if ln is not None:
        pd = ln.find('{%s}prstDash' % NS['a'])
        if pd is not None:
            val = pd.get('val', '')
            dash_ok = any(d in val for d in ('dash', 'sysDash', 'lgDash', 'dashDot', 'dot'))
    return (
        is_oval_like(s) and
        in_box(l, t, 10.5, 23.5, 6, 12) and
        l + w <= 23.5 and t + h <= 12 and
        w > h and
        in_wh(w, h, 11, 12.5, 4.5, 5.5) and
        fill_ok and
        line_rgb and is_light_orange(line_rgb) and
        lw and in_range(lw, 1.0, 1.5) and
        dash_ok
    )


# ── D2-14 +1 最内侧横向椭圆虚线环 ──
@rule('D2-14', +1, '中部最内侧横向椭圆虚线环')
def _(slide, prs):
    for s in slide.shapes:
        if inner_ellipse_ring_ok(s):
            l, t, w, h = shape_cm(s)
            lw = get_line_width_pt(s)
            return True, f'{s.name} {w:.1f}x{h:.1f}cm 最内侧浅橙色无填充虚线椭圆环 线宽{lw}pt'
    return False, '未找到符合距左10.5–23.5cm、距上6–12cm、宽11–12.5cm、高4.5–5.5cm、无填充、浅橙色均匀虚线、线宽1–1.5磅要求的最内侧横向椭圆环'


# ── D2-15 +3 三个椭圆环汇入上下箭头 ──
@rule('D2-15', +3, '三椭圆环汇入上下渐进式箭头相交处')
def _(slide, prs):
    """
    三个椭圆环一起汇入上方和下方中间与渐进式箭头相交处。
    检查（几何相交/接触，小容差）：
      1. 找到三个符合条件的椭圆环（按 D2-12/13/14 尺寸带、橙色系边线）。
      2. 找到顶部 upArrow 主体（D2-02）与两个渐进式 downArrow（D2-06 / D2-11 区域）。
      3. 三个椭圆环 + 三支箭头的水平中心 cx 严格共轴（最大偏差 ≤ TOL_X）。
      4. 每个椭圆环的水平中心 cx 严格落入 upArrow 的实际水平范围（不再向外扩容差），
         保证椭圆环轴线与上方箭头轴线对齐（真正"汇入"，非平移接近）。
      5. 每个椭圆环的顶部 y 与 upArrow 底边(=upArrow.tail) 接触/相交（间距 ≤ TOL_Y），
         即椭圆环顶部端点与上方渐进式箭头端点相交。
      6. 每个椭圆环的底部 (cx, bottom) 严格落入至少一个 downArrow 水平范围，
         且 y 与该 downArrow 顶边(=downArrow.tail) 接触/相交（间距 ≤ TOL_Y）。
    """
    TOL_X = 0.3   # cm，水平中心/端点水平容差（收紧至渲染/线宽级）
    TOL_Y = 0.35  # cm，接触/相交纵向容差（收紧至渲染/线宽级）

    # ── 找三个椭圆环 ──
    ovals = []
    for s in slide.shapes:
        if not is_oval_like(s):
            continue
        l, t, w, h = shape_cm(s)
        cx, cy = shape_center(s)
        line_rgb = get_line_rgb(s)
        if not (line_rgb and (is_light_orange(line_rgb) or is_orange(line_rgb) or is_deep_orange(line_rgb))):
            continue
        # 三层椭圆环均在中部水平带，宽 ≥ 11cm，高 ≥ 2cm
        if in_range(cx, 13, 21) and in_range(cy, 6.5, 12) and w >= 11 and h >= 2.0:
            ovals.append(s)

    if len(ovals) < 3:
        return False, f'仅找到{len(ovals)}个符合条件的椭圆环（需3个）'

    # ── 找顶部上箭头（上方渐进式箭头）──
    top_arrow = next(
        (s for s in slide.shapes
         if prst_geom(s) == 'upArrow' and
         in_range(shape_center(s)[0], 13, 22) and
         in_range(shape_center(s)[1], 0.1, 7)),
        None
    )
    if not top_arrow:
        return False, '未找到顶部上箭头主体（upArrow），无法判断汇入'

    # ── 找两个渐进式向下箭头（D2-06 与 D2-11 区域的 downArrow）──
    down_arrows = [
        s for s in slide.shapes
        if prst_geom(s) == 'downArrow' and
        in_range(shape_center(s)[0], 15, 19) and
        in_range(shape_center(s)[1], 6, 13)
    ]
    if len(down_arrows) < 2:
        return False, f'仅找到{len(down_arrows)}个渐进式向下箭头（需2个）'

    # ── (a) 水平中心整体对齐（共轴汇聚）──
    all_cx = (
        [shape_center(o)[0] for o in ovals] +
        [shape_center(top_arrow)[0]] +
        [shape_center(a)[0] for a in down_arrows]
    )
    spread = max(all_cx) - min(all_cx)
    if spread > TOL_X:
        return False, (
            f'椭圆环与上下渐进式箭头水平中心偏差{spread:.2f}cm > {TOL_X}cm，未共轴'
        )

    # ── (b) 椭圆环顶部与上方箭头相接/相交 ──
    #     · cx 必须严格落入 upArrow 实际水平范围（不向外扩容差）——真正的"汇入"要求
    #       椭圆环长轴中心线与箭头轴线重合，而非位于箭头旁边。
    #     · 顶部 y 与 upArrow 底边（=尾端）距离 ≤ TOL_Y，即椭圆环顶部端点与
    #       上方渐进式箭头端点几何接触/相交。
    ta_l, ta_t, ta_w, ta_h = shape_cm(top_arrow)
    ta_left, ta_right = ta_l, ta_l + ta_w
    ta_bottom = ta_t + ta_h  # upArrow 的底边（尾部所在边）
    for o in ovals:
        ol, ot, ow, oh = shape_cm(o)
        ocx = ol + ow / 2
        if not (ta_left <= ocx <= ta_right):
            return False, (
                f'椭圆环顶部 cx={ocx:.2f}cm 未严格落入 upArrow 水平范围'
                f'[{ta_left:.2f},{ta_right:.2f}]cm，未汇入箭头轴线'
            )
        if abs(ot - ta_bottom) > TOL_Y:
            return False, (
                f'椭圆环顶部 y={ot:.2f}cm 与 upArrow 底边 y={ta_bottom:.2f}cm '
                f'间距{abs(ot-ta_bottom):.2f}cm > {TOL_Y}cm，端点未接触/相交'
            )

    # ── (c) 椭圆环底部与两个 downArrow 之一相接/相交 ──
    #     · cx 必须严格落入某个 downArrow 实际水平范围；
    #     · 底部 y 与该 downArrow 顶边（=尾端）距离 ≤ TOL_Y，即椭圆环底部端点与
    #       下方渐进式箭头端点几何接触/相交。
    da_geo = []
    for a in down_arrows:
        al, at, aw, ah = shape_cm(a)
        da_geo.append((al, at, al + aw, at + ah))  # (left, top, right, bottom)
    for o in ovals:
        ol, ot, ow, oh = shape_cm(o)
        ocx = ol + ow / 2
        obottom = ot + oh
        touched = False
        for (dal, dat, dar, _dab) in da_geo:
            if not (dal <= ocx <= dar):
                continue
            if abs(obottom - dat) <= TOL_Y:
                touched = True
                break
        if not touched:
            return False, (
                f'椭圆环底部 (cx={ocx:.2f},y={obottom:.2f})cm 未与任一渐进式 '
                f'downArrow 的顶边端点接触/相交（容差{TOL_Y}cm）'
            )

    return True, (
        f'3个椭圆环共轴（cx偏差{spread:.2f}cm ≤ {TOL_X}cm），顶部与上方 upArrow 相接、'
        f'底部与2个渐进式 downArrow 相接（容差{TOL_Y}cm）'
    )


def _find_outer_ring(slide):
    return next((s for s in slide.shapes if outer_ellipse_ring_ok(s)), None)

def _left_side_connection_lines(slide, lo_l, hi_l, lo_t, hi_t):
    """在给定 bbox（用中心点判断）内收集所有直线/弧线/连接符，供左半辅助线相交检验。"""
    lines = []
    for s in slide.shapes:
        cx, cy = shape_center(s)
        if not (in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t)):
            continue
        if _is_office_line_or_connector(s):
            lines.append(s)
    return lines


def _find_left_aux_line(slide):
    """
    D2-16 专用：中部椭圆区域左半侧的弧形辅助连线。
    参考 _find_right_aux_line 的实现，除位置/颜色/线宽外，
    再验证其与"左上侧连接线"（D2-44 区域：距左 9–12.5cm、距上 5.2–7.5cm）
    以及"左下侧连接线"（D2-45 区域：距左 9.5–13cm、距上 11.0–13.5cm）
    的边界均相交，作为"与左上、左下侧连接线相交"的实质校验。
    """
    top_lines = _left_side_connection_lines(slide, 9.0, 12.5, 5.2, 7.5)
    bottom_lines = _left_side_connection_lines(slide, 9.5, 13.0, 11.0, 13.5)

    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        # 细则位置：距左 9–12.6cm、距上 8.5–10.5cm（以左上角判断，与 D2-18 对称）
        if not (in_range(l, 9.0, 12.6) and in_range(t, 8.5, 10.5)):
            continue
        # 弧形/线条/连接符/自由曲线均可（与右侧辅助线一致）
        if not _is_office_line_or_connector(s):
            continue
        line_rgb = get_line_rgb(s)
        if not (line_rgb and is_light_orange(line_rgb)):
            continue
        lw = get_line_width_pt(s)
        if not (lw and in_range(lw, 1.0, 1.5)):
            continue
        # 排除大面积填充块（宽度和高度都超过 5cm 的不是辅助线）
        if w > 5 and h > 5:
            continue
        intersects_top = any(_bounds_intersect(s, line) for line in top_lines)
        intersects_bottom = any(_bounds_intersect(s, line) for line in bottom_lines)
        if not (intersects_top and intersects_bottom):
            continue
        return s
    return None

def _is_office_line_or_connector(shape):
    """Office/WPS 中的直线、弧线、连接符或自由曲线对象。"""
    geom = prst_geom(shape)
    tag = shape._element.tag.split('}')[-1]
    if tag == 'cxnSp':
        return True
    if geom in ('line', 'arc', 'curvedConnector2', 'curvedConnector3', 'curvedConnector4', 'curvedConnector5'):
        return True
    if shape._element.find('.//{%s}custGeom' % NS['a']) is not None:
        return True
    return is_thin_line(shape)


def _bounds_intersect(a, b, tol=0.25):
    al, at, aw, ah = shape_cm(a)
    bl, bt, bw, bh = shape_cm(b)
    ar, ab = al + aw, at + ah
    br, bb = bl + bw, bt + bh
    return not (ar + tol < bl or br + tol < al or ab + tol < bt or bb + tol < at)


def _right_side_connection_lines(slide, lo_l, hi_l, lo_t, hi_t):
    lines = []
    for s in slide.shapes:
        cx, cy = shape_center(s)
        if not (in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t)):
            continue
        if _is_office_line_or_connector(s):
            lines.append(s)
    return lines


def _find_right_aux_line(slide):
    """
    D2-18：第1页中部右半椭圆辅助线。
    仅检查细则要求：右半侧区域、弧形辅助连线、浅橙色、线宽1–1.5磅、
    与右上和右下侧连接线相交。
    """
    top_lines = _right_side_connection_lines(slide, 21.0, 24.1, 5.2, 7.5)
    bottom_lines = _right_side_connection_lines(slide, 21.5, 24.5, 11.0, 13.5)

    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        if not (in_range(l, 21.5, 26.0) and in_range(t, 8.5, 10.5)):
            continue
        if not _is_office_line_or_connector(s):
            continue
        line_rgb = get_line_rgb(s)
        if not (line_rgb and is_light_orange(line_rgb)):
            continue
        lw = get_line_width_pt(s)
        if not (lw and in_range(lw, 1.0, 1.5)):
            continue
        intersects_top = any(_bounds_intersect(s, line) for line in top_lines)
        intersects_bottom = any(_bounds_intersect(s, line) for line in bottom_lines)
        if not (intersects_top and intersects_bottom):
            continue
        return s
    return None


# ── D2-16 +1 左半椭圆辅助线 ──
@rule('D2-16', +1, '中部左半椭圆辅助线')
def _(slide, prs):
    s = _find_left_aux_line(slide)
    if s:
        l, t, w, h = shape_cm(s)
        lw = get_line_width_pt(s)
        return True, f'{s.name} 左半椭圆辅助连线 位置({l:.1f},{t:.1f}) 尺寸{w:.1f}x{h:.1f}cm 线宽{lw}pt'
    return False, ('未找到左半椭圆辅助连线：需同时满足 距左9–12.6/距上8.5–10.5cm、'
                   '弧线/连接符/自由曲线、浅橙色、线宽1–1.5磅，'
                   '且与左上侧连接线(9–12.5×5.2–7.5cm)和左下侧连接线(9.5–13×11–13.5cm)边界相交')


# ── D2-17 +1 最外层椭圆左侧圆点 ──
@rule('D2-17', +1, '最外层椭圆左侧中心圆点')
def _(slide, prs):
    """
    细则：最外层横向椭圆环左侧中心有一个宽高 0.4–0.5cm 的圆点，与左半椭圆辅助线相交。
    检查点（仅按细则约束，细则未提及的不做限制）：
      1. 形状为圆形（is_oval_like）
      2. 宽高均在 0.4–0.5cm
      3. 位于最外层椭圆环左侧中心：
           以最外层椭圆的左边缘 X、垂直中心 Y 为参考；
           圆点中心与参考点的偏差：X ≤ 1.0cm，Y ≤ 0.8cm。
           若最外层椭圆未找到，此项直接失败。
      4. 与左半椭圆辅助线相交：
           圆点中心 X 与辅助线中心 X 偏差 ≤ 1.0cm，
           圆点中心 Y 与辅助线中心 Y 偏差 ≤ 0.8cm。
           若辅助线未找到，此项直接失败。
    """
    outer = _find_outer_ring(slide)
    if not outer:
        return False, '未找到最外层横向椭圆环，无法确认圆点位于其左侧中心'

    aux_line = _find_left_aux_line(slide)
    if not aux_line:
        return False, '未找到左半椭圆辅助线（D2-16），无法验证圆点与辅助线相交'

    ol, ot, ow, oh = shape_cm(outer)
    ref_x = ol            # 椭圆左边缘 X
    ref_y = ot + oh / 2   # 椭圆垂直中心 Y
    ax, ay = shape_center(aux_line)

    for s in slide.shapes:
        # ── 1. 圆形 ──
        if not is_oval_like(s):
            continue
        # ── 2. 宽高 0.4–0.5cm ──
        l2, t2, w2, h2 = shape_cm(s)
        if not (in_range(w2, 0.4, 0.5) and in_range(h2, 0.4, 0.5)):
            continue
        cx, cy = shape_center(s)
        # ── 3. 位于最外层椭圆左侧中心附近 ──
        if not (abs(cx - ref_x) <= 1.0 and abs(cy - ref_y) <= 0.8):
            continue
        # ── 4. 与左半椭圆辅助线相交 ──
        if not (abs(cx - ax) <= 1.0 and abs(cy - ay) <= 0.8):
            continue
        return True, (
            f'{s.name} 左侧中心圆点 ({l2:.2f},{t2:.2f}) '
            f'{w2:.2f}x{h2:.2f}cm 与左半椭圆辅助线相交'
        )

    return False, (
        f'未找到宽高0.4–0.5cm、位于最外层椭圆左侧中心'
        f'（参考点 x≈{ref_x:.1f} y≈{ref_y:.1f}cm，偏差≤1.0/0.8cm）'
        f'且与左半椭圆辅助线（中心 x≈{ax:.1f} y≈{ay:.1f}cm）相交的圆形圆点'
    )


# ── D2-18 +1 右半椭圆辅助线 ──
@rule('D2-18', +1, '中部右半椭圆辅助线')
def _(slide, prs):
    """
    细则：位于中部椭圆区域右半侧，距左21.5–26cm、距上8.5–10.5cm范围内，
    为一条弧形辅助连线，颜色为浅橙色，线宽1–1.5磅，与右上、右下侧连接线相交。
    检查点（仅按细则约束）：
      1. 位置：left∈[21.5,26.0]cm，top∈[8.5,10.5]cm
      2. 为弧形辅助连线（Office/WPS线条或连接符对象）
      3. 颜色为浅橙色
      4. 线宽1–1.5磅
      5. 与右上侧连接线（D2-46区域）相交
      6. 与右下侧连接线（D2-47区域）相交
    """
    s = _find_right_aux_line(slide)
    if s:
        l, t, w, h = shape_cm(s)
        lw = get_line_width_pt(s)
        return True, (
            f'{s.name} 右半椭圆辅助连线 位置({l:.2f},{t:.2f}) '
            f'尺寸{w:.2f}x{h:.2f}cm 浅橙色 线宽{lw}pt 与右上右下连接线相交'
        )
    return False, (
        '未找到符合以下全部条件的右半椭圆辅助连线：'
        'left∈[21.5,26.0]cm、top∈[8.5,10.5]cm、'
        '浅橙色、线宽1–1.5磅、与右上侧连接线相交且与右下侧连接线相交'
    )


# ── D2-19 +1 最外层椭圆右侧圆点 ──
@rule('D2-19', +1, '最外层椭圆右侧中心圆点')
def _(slide, prs):
    """
    细则：最外层横向椭圆环右侧中心有一个宽高 0.4–0.5cm 的圆点，与右半椭圆辅助线相交。
    检查点（仅按细则约束，细则未提及的不做限制）：
      1. 形状为圆形（is_oval_like）
      2. 宽高均在 0.4–0.5cm
      3. 位于最外层椭圆环右侧中心：
           以最外层椭圆的右边缘 X、垂直中心 Y 为参考；
           圆点中心与参考点的偏差：X ≤ 1.0cm，Y ≤ 0.8cm。
           若最外层椭圆未找到，此项直接失败。
      4. 与右半椭圆辅助线相交：
           使用 Office/WPS 图形边界框相交判断，允许 0.25cm 误差。
           若右半椭圆辅助线未找到，此项直接失败。
    """
    outer = _find_outer_ring(slide)
    if not outer:
        return False, '未找到最外层横向椭圆环，无法确认圆点位于其右侧中心'

    aux_line = _find_right_aux_line(slide)
    if not aux_line:
        return False, '未找到右半椭圆辅助线（D2-18），无法验证圆点与辅助线相交'

    ol, ot, ow, oh = shape_cm(outer)
    ref_x = ol + ow       # 椭圆右边缘 X
    ref_y = ot + oh / 2   # 椭圆垂直中心 Y

    for s in slide.shapes:
        # ── 1. 圆形 ──
        if not is_oval_like(s):
            continue
        # ── 2. 宽高 0.4–0.5cm ──
        l2, t2, w2, h2 = shape_cm(s)
        if not (in_range(w2, 0.4, 0.5) and in_range(h2, 0.4, 0.5)):
            continue
        cx, cy = shape_center(s)
        # ── 3. 位于最外层椭圆右侧中心附近 ──
        if not (abs(cx - ref_x) <= 1.0 and abs(cy - ref_y) <= 0.8):
            continue
        # ── 4. 与右半椭圆辅助线相交 ──
        if not _bounds_intersect(s, aux_line, tol=0.25):
            continue
        return True, (
            f'{s.name} 右侧中心圆点 ({l2:.2f},{t2:.2f}) '
            f'{w2:.2f}x{h2:.2f}cm 与右半椭圆辅助线相交'
        )

    return False, (
        f'未找到宽高0.4–0.5cm、位于最外层椭圆右侧中心'
        f'（参考点 x≈{ref_x:.1f} y≈{ref_y:.1f}cm，偏差≤1.0/0.8cm）'
        f'且与右半椭圆辅助线相交的圆形圆点'
    )


def _round_rect_adj_cm(shape):
    """
    读取圆角矩形的圆角半径（cm）。
    Office/WPS 将圆角大小存储为 adjVal，取值 0–100000（对应 0%–100% 的半高）。
    实际圆角半径 = adjVal/100000 × min(width, height) / 2。
    若读取失败则返回 None。
    """
    try:
        prstGeom = shape._element.find('.//{%s}prstGeom' % NS['a'])
        if prstGeom is None:
            return None
        avLst = prstGeom.find('{%s}avLst' % NS['a'])
        if avLst is None:
            return None
        gd = avLst.find('{%s}gd' % NS['a'])
        if gd is None:
            return None
        fmla = gd.get('fmla', '')
        # fmla 形如 "val 16667"
        parts = fmla.strip().split()
        if len(parts) == 2 and parts[0] == 'val':
            adj = int(parts[1])
            l, t, w, h = shape_cm(shape)
            radius = adj / 100000.0 * min(w, h) / 2.0
            return radius
    except Exception:
        pass
    return None


def _find_impl_frame(slide):
    """
    细则：位于距左11.5–22.5cm、距上7.0–10.9cm范围内，
    宽9.5–10.2cm，高2.8–3.3cm，形状为圆角矩形，
    填充为浅橙色，边线为深橙色单实线，线宽2–3磅，
    圆角半径0.35–0.7cm。
    """
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        # 位置
        if not in_box(l, t, 11.5, 22.5, 7.0, 10.9):
            continue
        # 尺寸
        if not in_wh(w, h, 9.5, 10.2, 2.8, 3.3):
            continue
        # 形状为圆角矩形
        if prst_geom(s) != 'roundRect':
            continue
        # 填充为浅橙色
        rgb, _ = get_fill_rgb(s)
        pale_orange = bool(rgb and rgb[0] >= 230 and rgb[1] >= 190 and rgb[2] <= 220 and rgb[0] > rgb[2])
        if not pale_orange:
            continue
        # 边线为深橙色单实线，线宽2–3磅，无虚线
        line_rgb = get_line_rgb(s)
        lw = get_line_width_pt(s)
        if not (line_rgb and is_deep_orange(line_rgb)):
            continue
        if not (lw and in_range(lw, 2, 3)):
            continue
        if has_dash(s):
            continue
        # 圆角半径 0.35–0.7cm
        radius = _round_rect_adj_cm(s)
        if radius is not None and not in_range(radius, 0.35, 0.7):
            continue
        return s
    return None


# ── D2-20 +1 中部实施框外层圆角矩形 ──
@rule('D2-20', +1, '中部实施框外层圆角矩形')
def _(slide, prs):
    """
    细则：位于距左11.5–22.5cm、距上7.0–10.9cm范围内，
    宽9.5–10.2cm，高2.8–3.3cm，形状为圆角矩形，
    填充为浅橙色，边线为深橙色单实线，线宽2–3磅，圆角半径0.35–0.7cm。
    """
    s = _find_impl_frame(slide)
    if s:
        l, t, w, h = shape_cm(s)
        lw = get_line_width_pt(s)
        radius = _round_rect_adj_cm(s)
        return True, (
            f'{s.name} roundRect ({l:.2f},{t:.2f}) {w:.2f}x{h:.2f}cm '
            f'浅橙填充 深橙单实线边框{lw}pt '
            f'圆角半径{"%.2f" % radius if radius is not None else "未读取"}cm'
        )
    return False, (
        '未找到符合以下全部条件的中部实施框圆角矩形：'
        'left∈[11.5,22.5]cm、top∈[7.0,10.9]cm、'
        '宽9.5–10.2cm、高2.8–3.3cm、圆角矩形、'
        '浅橙色填充、深橙色单实线边线、线宽2–3磅、圆角半径0.35–0.7cm'
    )


# ── D2-21 +1 实施框标题文本"实施枢纽" ──
@rule('D2-21', +1, '实施框标题文本"实施枢纽"')
def _(slide, prs):
    """
    细则：位于中部实施框上半部，文本为"实施枢纽"，
    字体为黑体、微软雅黑或相近字体，字号21–23磅，加粗，
    颜色为深橙色，水平居中。
    检查点（仅按细则约束）：
      1. 文本精确等于"实施枢纽"
      2. 位于中部实施框上半部（实施框 top 到 top+height×0.55 范围内）
         若实施框未找到则退回绝对区域 left∈[9.0,24.5]、top∈[7.0,9.5]
      3. 字号 21–23 磅
      4. 加粗
      5. 颜色为深橙色
      6. 字体为黑体/微软雅黑或相近中文无衬线字体
      7. 水平居中
    """
    s = find_shape_with_text(slide, '实施枢纽')
    if not s:
        return False, '未找到文本精确等于"实施枢纽"的形状'

    # ── 2. 位置：中部实施框上半部 ──
    impl = _find_impl_frame(slide)
    l, t, w, h = shape_cm(s)
    if impl:
        il, it, iw, ih = shape_cm(impl)
        upper_half_hi = it + ih * 0.55
        if not (t >= it and t <= upper_half_hi):
            return False, (
                f'"实施枢纽"top={t:.2f}cm 不在实施框上半部'
                f'（实施框 top={it:.2f}cm，上半部上限={upper_half_hi:.2f}cm）'
            )
    else:
        if not in_box(l, t, 9.0, 24.5, 7.0, 9.5):
            return False, f'"实施枢纽"位置({l:.2f},{t:.2f})不在中部实施框上半部区域'

    # ── 3–7. 字体属性 ──
    all_runs = _font_info_all_runs(s)
    if not all_runs:
        fi = font_info(s)
        if fi:
            all_runs = [fi]
    if not all_runs:
        return False, '"实施枢纽"无法读取字体信息'

    sz_ok    = any(r[0] and in_range(r[0], 21, 23) for r in all_runs)
    bold_ok  = any(_is_bold_or_semibold(r[1], r[3]) for r in all_runs)
    color_ok = any(r[2] and (is_deep_orange(r[2]) or is_orange(r[2])) for r in all_runs)
    font_ok  = any(is_cjk_sans(r[3]) for r in all_runs)
    center_ok = text_is_centered(s)

    if not sz_ok:
        return False, f'"实施枢纽"字号{[r[0] for r in all_runs]}不在21–23磅范围'
    if not bold_ok:
        return False, '"实施枢纽"未加粗'
    if not color_ok:
        return False, f'"实施枢纽"颜色{[r[2] for r in all_runs]}非深橙色'
    if not font_ok:
        return False, f'"实施枢纽"字体{[r[3] for r in all_runs]}非黑体/微软雅黑或相近字体'
    if not center_ok:
        return False, '"实施枢纽"未水平居中'

    return True, (
        f'{s.name} "实施枢纽" 位置({l:.2f},{t:.2f}) '
        f'字号{[r[0] for r in all_runs if r[0]]}pt 加粗 深橙色 居中'
    )


# ── D2-22 +1 实施框标题下装饰线与圆点 ──
@rule('D2-22', +1, '实施框标题下装饰线与圆点')
def _(slide, prs):
    """
    细则：位于"实施枢纽"文本下方0.35cm–0.5cm处，
    为一条浅橙色水平单实线，长度7cm–8cm，线宽0.75–1.25磅；
    装饰线中间有一个深橙色小圆点，直径0.12cm–0.22cm。
    检查点（仅按细则约束）：
      1. 以"实施枢纽"文本底边为参考，装饰线中心Y在其下方0.35–0.5cm
      2. 装饰线为水平单实线，长度7–8cm，浅橙色，线宽0.75–1.25磅
      3. 小圆点为圆形，宽高均0.12–0.22cm，深橙色
      4. 小圆点位于装饰线中间：圆点中心X与线中心X偏差≤0.15cm，中心Y与线中心Y偏差≤0.15cm
    """
    hub = find_shape_with_text(slide, '实施枢纽')
    if not hub:
        return False, '未找到"实施枢纽"文本作为装饰线参考位置'

    hub_l, hub_t, hub_w, hub_h = shape_cm(hub)
    hub_bottom = hub_t + hub_h
    lines = []
    dots = []

    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        cx, cy = shape_center(s)
        lr = get_line_rgb(s)
        lw = get_line_width_pt(s)
        # 装饰线：位于文本下方0.35–0.5cm，浅橙色水平单实线，长度7–8cm，线宽0.75–1.25磅
        if (is_horizontal_line(s, 7, 8, 0.1) and
                lr and is_light_orange(lr) and
                lw and in_range(lw, 0.75, 1.25) and
                not has_dash(s) and
                in_range(cy - hub_bottom, 0.35, 0.5)):
            lines.append(s)
        # 深橙色小圆点，直径0.12–0.22cm
        if is_oval_like(s) and in_range(w, 0.12, 0.22) and in_range(h, 0.12, 0.22):
            rgb, _ = get_fill_rgb(s)
            if rgb and is_deep_orange(rgb):
                dots.append(s)

    if not lines:
        return False, '未找到位于"实施枢纽"文本下方0.35–0.5cm、长度7–8cm、浅橙色、线宽0.75–1.25磅的水平单实线'
    if not dots:
        return False, '未找到直径0.12–0.22cm的深橙色小圆点'

    for line in lines:
        lx, ly = line_center(line)
        for dot in dots:
            dx, dy = shape_center(dot)
            if abs(dx - lx) <= 0.15 and abs(dy - ly) <= 0.15:
                ll, lt, lw_cm, lh_cm = shape_cm(line)
                dl, dt, dw, dh = shape_cm(dot)
                return True, (
                    f'{line.name} 位于"实施枢纽"下方{ly-hub_bottom:.2f}cm，'
                    f'长度{lw_cm:.2f}cm 浅橙色水平单实线；'
                    f'{dot.name} 深橙色小圆点直径约{(dw+dh)/2:.2f}cm 位于装饰线中间'
                )

    return False, (
        f'找到装饰线{len(lines)}条、深橙色小圆点{len(dots)}个，'
        f'但未满足小圆点位于装饰线中间（中心偏差≤0.15cm）'
    )


def _normalize_content_text(text):
    return re.sub(r'\s+', '', text).replace('·', '·')


def _is_dark_or_black(rgb):
    # 非主题色（深色/黑色）：只要视觉上属于深色即通过，不要求接近纯黑或具体主题色。
    if rgb is None:
        return False
    # 允许深灰、深蓝灰等各类"深色"，不强制通道近似相等
    return max(rgb) <= 128


# ── D2-23 +1 实施框内容文本 ──
@rule('D2-23', +1, '实施框内容文本"课程内容 · 技术平台 · 活动组织"')
def _(slide, prs):
    """
    细则：位于中部实施框下半部，文本为"课程内容 · 技术平台 · 活动组织"，
    字体为黑体、微软雅黑或相近字体，字号13–15磅，加粗或半加粗，
    颜色为深色或黑色，水平居中。
    检查点（仅按细则约束）：
      1. 文本规范化后精确等于"课程内容·技术平台·活动组织"
      2. 位于中部实施框下半部（实施框 top+height×0.45 到 bottom）
      3. 字体为黑体/微软雅黑或相近中文无衬线字体
      4. 字号13–15磅
      5. 加粗或半加粗
      6. 颜色为深色或黑色
      7. 水平居中
    """
    impl = _find_impl_frame(slide)
    if not impl:
        return False, '未找到中部实施框，无法确认内容文本位于实施框下半部'

    il, it, iw, ih = shape_cm(impl)
    lower_top = it + ih * 0.45
    lower_bottom = it + ih
    expected = '课程内容·技术平台·活动组织'

    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        raw_text = shape_text(s).strip()
        if _normalize_content_text(raw_text) != expected:
            continue

        l, t, w, h = shape_cm(s)
        if not (t >= lower_top and t <= lower_bottom):
            return False, (
                f'内容文本top={t:.2f}cm 不在实施框下半部'
                f'（下半部范围 {lower_top:.2f}–{lower_bottom:.2f}cm）'
            )

        all_runs = _font_info_all_runs(s)
        if not all_runs:
            fi = font_info(s)
            if fi:
                all_runs = [fi]
        if not all_runs:
            return False, '内容文本无法读取字体信息'

        font_ok = any(is_cjk_sans(r[3]) for r in all_runs)
        sz_ok = any(r[0] and in_range(r[0], 13, 15) for r in all_runs)
        bold_ok = any(_is_bold_or_semibold(r[1], r[3]) for r in all_runs)
        color_ok = any(r[2] and _is_dark_or_black(r[2]) for r in all_runs)
        center_ok = text_is_centered(s)

        if not font_ok:
            return False, f'内容文本字体{[r[3] for r in all_runs]}非黑体/微软雅黑或相近字体'
        if not sz_ok:
            return False, f'内容文本字号{[r[0] for r in all_runs]}不在13–15磅范围'
        if not bold_ok:
            return False, '内容文本未加粗或半加粗'
        if not color_ok:
            return False, f'内容文本颜色{[r[2] for r in all_runs]}非深色或黑色'
        if not center_ok:
            return False, '内容文本未水平居中'

        return True, (
            f'{s.name} 内容文本位于实施框下半部 位置({l:.2f},{t:.2f}) '
            f'字号{[r[0] for r in all_runs if r[0]]}pt 加粗/半加粗 深色/黑色 居中'
        )

    return False, '未找到文本精确为"课程内容 · 技术平台 · 活动组织"的实施框内容文本'


# ─────────────────── 卡片规则生成器 ───────────────────
def broad_region_for(title):
    return None  # 宽松 broad 匹配已移除，只用严格区域

def card_frame_checker(slide, region, side, title):
    lo_l, hi_l, lo_t, hi_t, lo_w, hi_w, lo_h, hi_h = region
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        if not in_box(l, t, lo_l, hi_l, lo_t, hi_t): continue
        if not in_wh(w, h, lo_w, hi_w, lo_h, hi_h): continue
        if prst_geom(s) not in ('parallelogram', 'roundRect'): continue
        rgb, _ = get_fill_rgb(s)
        line_rgb = get_line_rgb(s)
        lw = get_line_width_pt(s)
        if (rgb is None or is_white_or_near(rgb)) and line_rgb and is_orange(line_rgb) and lw and in_range(lw, 1.5, 2.2):
            return True, f'{title}外框 {s.name} ({l:.1f},{t:.1f}) {w:.1f}x{h:.1f}cm 白底橙边'
    return False, f'未找到{title}外框'

def icon_checker(slide, region, title):
    lo_l, hi_l, lo_t, hi_t = region
    circle = None
    internal = False
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        if not in_box(l, t, lo_l, hi_l, lo_t, hi_t): continue
        if l + w > hi_l or t + h > hi_t: continue
        rgb, _ = get_fill_rgb(s)
        if is_oval_like(s) and in_range(w, 1.3, 1.6) and in_range(h, 1.3, 1.6) and rgb and is_orange(rgb):
            circle = s
    if circle:
        cx, cy = shape_center(circle)
        for s in slide.shapes:
            if s == circle: continue
            x, y = shape_center(s)
            if abs(x-cx) <= 0.75 and abs(y-cy) <= 0.75:
                rgb, _ = get_fill_rgb(s)
                trgb = text_color_rgb(s)
                if (rgb and is_white_or_near(rgb)) or (trgb and is_white_or_near(trgb)):
                    internal = True
        if internal:
            return True, f'{title}图标含橙色圆形 {circle.name} 和白色内部可编辑元素'
    return False, f'未找到{title}带白色内部图形的橙色圆形图标'

def _card_frame_passed(slide, rid_frame):
    for rid, score, title, fn in RULES:
        if rid == rid_frame:
            ok, _ = fn(slide, None)
            return ok
    return False

def title_checker(slide, text, region, frame_rid=None):
    if frame_rid and not _card_frame_passed(slide, frame_rid):
        return False, f'卡片外框({frame_rid})未通过，标题不计分'
    lo_l, hi_l, lo_t, hi_t = region
    s = find_shape_with_text(slide, text)
    if not s: return False, f'未找到"{text}"'
    l, t, w, h = shape_cm(s)
    if not in_box(l, t, lo_l, hi_l, lo_t, hi_t): return False, f'"{text}"位置({l:.1f},{t:.1f})不在范围'
    fi = font_info(s)
    if not fi: return False, '无字体信息'
    sz, bold, color, name, align = fi
    if sz and in_range(sz, 16, 19.5) and bold and is_cjk_sans(name):
        return True, f'{text} 字号{sz}pt 加粗'
    return False, f'{text}样式不满足 sz={sz} bold={bold} font={name}'

def sep_line_dot_checker(slide, region, title, frame_rid=None):
    if frame_rid and not _card_frame_passed(slide, frame_rid):
        return False, f'卡片外框({frame_rid})未通过，分隔线不计分'
    lo_l, hi_l, lo_t, hi_t = region
    line_found = False; dot_found = False
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        if not in_box(l, t, lo_l, hi_l, lo_t, hi_t): continue
        line_rgb = get_line_rgb(s)
        if in_range(w,3.5,4.5) and h <= 0.2 and line_rgb and (is_light_orange(line_rgb) or is_orange(line_rgb)):
            lw = get_line_width_pt(s)
            if lw and in_range(lw, 0.75, 1.25): line_found = True
        if is_oval_like(s) and in_range(w,0.12,0.35) and in_range(h,0.12,0.35):
            rgb, _ = get_fill_rgb(s)
            if rgb and is_orange(rgb): dot_found = True
    if line_found and dot_found:
        return True, f'{title}分隔线与圆点存在'
    return False, f'{title}分隔线={line_found} 圆点={dot_found}'

def list_text_checker(slide, keywords, region, size_range, title, frame_rid=None):
    if frame_rid and not _card_frame_passed(slide, frame_rid):
        return False, f'卡片外框({frame_rid})未通过，三行文本不计分'
    lo_l, hi_l, lo_t, hi_t = region
    for s in slide.shapes:
        if not s.has_text_frame: continue
        txt = shape_text(s)
        if not all(k in txt for k in keywords): continue
        l, t, w, h = shape_cm(s)
        if not in_box(l, t, lo_l, hi_l, lo_t, hi_t): continue
        fi = font_info(s)
        if fi:
            sz, bold, color, name, align = fi
            if sz and in_range(sz, size_range[0], size_range[1]) and is_cjk_sans(name) and color and not is_white_or_near(color):
                return True, f'{title}三行文本包含 {"、".join(keywords)}，字号{sz}pt'
    return False, f'未找到{title}三行文本'


# ── 四个侧边卡片 ──
@rule('D2-24', +3, '左上资源整合卡片外框')
def _(slide, prs):
    """
    细则：位于距左4–10.5cm、距上3.1–7.8cm范围内，宽4–5cm，高3.8–4.6cm，
    整体为左侧较宽、右侧略窄的斜边圆角四边形（parallelogram），
    填充为白色，边线为橙色单实线，线宽1.5–2.2磅。
    检查点（仅按细则约束）：
      1. 位置：left∈[4,10.5]cm，top∈[3.1,7.8]cm
      2. 尺寸：宽4–5cm，高3.8–4.6cm
      3. 形状为斜边圆角四边形（parallelogram），且**未**水平翻转 → 左侧较宽、右侧略窄
      4. 线条转角为圆角（<a:ln> 无 miter/bevel）→ 圆角表现
      5. 填充为白色（白色或接近白色）
      6. 边线为橙色单实线，线宽1.5–2.2磅
    """
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        if not in_box(l, t, 4, 10.5, 3.1, 7.8):
            continue
        if not in_wh(w, h, 4, 5, 3.8, 4.6):
            continue
        if prst_geom(s) != 'parallelogram':
            continue
        # 方向：左宽右窄 ⇒ parallelogram 未水平翻转
        if has_flip_h(s):
            continue
        # 圆角表现：线条 join 非 miter/bevel（PPT 默认或显式 round）
        if not has_round_corners(s):
            continue
        rgb, _ = get_fill_rgb(s)
        if not (rgb is None or is_white_or_near(rgb)):
            continue
        line_rgb = get_line_rgb(s)
        lw = get_line_width_pt(s)
        if not (line_rgb and is_orange(line_rgb)):
            continue
        if not (lw and in_range(lw, 1.5, 2.2)):
            continue
        if has_dash(s):
            continue
        return True, (
            f'{s.name} 左上资源整合卡片外框 ({l:.2f},{t:.2f}) '
            f'{w:.2f}x{h:.2f}cm parallelogram(未翻转→左宽右窄) 圆角 '
            f'白色填充 橙色单实线 线宽{lw:.2f}pt'
        )
    return False, (
        '未找到符合以下全部条件的左上资源整合卡片外框：'
        'left∈[4,10.5]cm、top∈[3.1,7.8]cm、宽4–5cm、高3.8–4.6cm、'
        'parallelogram(左宽右窄，未水平翻转)、圆角线条转角、白色填充、'
        '橙色单实线边线、线宽1.5–2.2磅'
    )

def _shape_inside_circle(shape, circle, margin=0.1):
    cl, ct, cw, ch = shape_cm(circle)
    l, t, w, h = shape_cm(shape)
    return (
        l >= cl - margin and t >= ct - margin and
        l + w <= cl + cw + margin and t + h <= ct + ch + margin
    )


def _is_white_editable_icon_part(shape):
    fill_rgb, _ = get_fill_rgb(shape)
    line_rgb = get_line_rgb(shape)
    text_rgb = text_color_rgb(shape)
    return (
        (fill_rgb and is_white_or_near(fill_rgb)) or
        (line_rgb and is_white_or_near(line_rgb)) or
        (text_rgb and is_white_or_near(text_rgb))
    )


def _is_open_book_icon(circle, white_parts):
    """判断圆内白色部件组是否构成"打开书本"几何图标。
    "打开的书本"的关键部件（细则的核心视觉特征）：
      · 书脊 = 位于图标水平中线附近的一根竖直细长部件（贯穿垂直方向的中缝），
              或一个跨中线的极窄部件（宽≪高 或 宽 ≤ 图标宽 × 20%）；
      · 书页 = 左右两侧对称展开的书页轮廓，至少一侧为曲线/自由几何/圆角矩形
              等非狭长闭合形状（不是纯直线/单段连接符）。
    通过判定（择一命中即可）：
      (A) 具备"书脊 + 至少两侧各一件书页部件（其中至少一侧为书页几何）"；
      (B) 单一跨中线的 custGeom/曲线自由形状（宽 ≥ 图标直径 45%、高 ≥ 图标直径 25%
          且宽/高比 ≤ 3.0）—— 整张打开书本作为一个自由几何绘制时，
          必须有足够高度体现两页展开轮廓，而不是一根扁细横线。
    否则返回 False，避免"任意 ≥2 个白色线条/形状"通过。"""
    if len(white_parts) < 2:
        return False, f'白色部件仅{len(white_parts)}个'

    cl, ct, cw, ch = shape_cm(circle)
    ccx = cl + cw / 2.0
    ccy = ct + ch / 2.0

    def _bounds(p):
        return shape_cm(p)  # (l, t, w, h)

    def _is_freeform_like(p):
        # custGeom / freeform：真正用于书页/整体书本的不规则轮廓
        if p._element.find('.//{%s}custGeom' % NS['a']) is not None:
            return True
        return prst_geom(p) == 'freeform'

    def _is_page_like(p):
        """书页几何：曲线/自由几何/圆角矩形/括号/书本相关预设 + 非狭长比例。"""
        if _is_freeform_like(p):
            return True
        pg = prst_geom(p)
        if pg.startswith('curved') or pg.startswith('wave') or pg in (
            'arc', 'foldedCorner', 'plaque', 'book', 'bookmark',
            'leftBrace', 'rightBrace', 'leftBracket', 'rightBracket', 'bracePair',
            'chevron', 'homePlate', 'roundRect', 'round1Rect', 'round2SameRect',
            'round2DiagRect',
        ):
            # 需附加"非狭长"：排除极扁的一根竖线冒充书页
            _pl, _ppt, pw, ph = _bounds(p)
            if pw > 0.05 and ph > 0.05:
                ratio = max(pw, ph) / max(min(pw, ph), 0.001)
                return ratio <= 6
            return False
        return False

    def _is_spine_like(p):
        """书脊：位于图标中线附近、竖直细长的部件（或跨中线极窄部件）。"""
        pl, pt2, pw, ph = _bounds(p)
        pcx = pl + pw / 2.0
        pcy = pt2 + ph / 2.0
        # 中线附近（水平中心距≤图标宽15%，垂直中心近图标中央）
        near_axis = abs(pcx - ccx) <= max(0.10, cw * 0.15)
        vert_mid = abs(pcy - ccy) <= max(0.30, ch * 0.35)
        if not (near_axis and vert_mid):
            return False
        # 竖直细长：高 ≥ 图标直径 35%、宽 ≤ 图标直径 22%、且高 > 宽
        if ph >= ch * 0.35 and pw <= cw * 0.22 and ph > pw:
            return True
        return False

    left_parts, right_parts, span_parts = [], [], []
    for p in white_parts:
        pl, _pt2, pw, _ph = _bounds(p)
        pcx = pl + pw / 2.0
        pr = pl + pw
        if pl < ccx - 0.05 and pr > ccx + 0.05:
            span_parts.append(p)
        elif pcx < ccx:
            left_parts.append(p)
        else:
            right_parts.append(p)

    # (A) 书脊 + 两页对称
    spine_parts = [p for p in white_parts if _is_spine_like(p)]
    if spine_parts and left_parts and right_parts:
        left_page = any(_is_page_like(p) for p in left_parts)
        right_page = any(_is_page_like(p) for p in right_parts)
        if left_page or right_page:
            return True, (
                f'书脊+两页 (脊{len(spine_parts)}件、左{len(left_parts)}/右{len(right_parts)}，'
                f'至少一侧为书页几何：左={left_page},右={right_page})'
            )

    # (B) 单一跨中线自由几何：整张打开书本作为一个 custGeom/freeform
    for p in span_parts:
        if not _is_freeform_like(p):
            continue
        _pl, _pt2, pw, ph = _bounds(p)
        if pw >= cw * 0.45 and ph >= ch * 0.25:
            ratio = pw / max(ph, 0.001)
            if ratio <= 3.0:
                return True, (
                    f'跨中线自由几何书本 {pw:.2f}x{ph:.2f}cm '
                    f'(宽≥{cw*0.45:.2f}、高≥{ch*0.25:.2f}、宽高比{ratio:.2f}≤3.0)'
                )

    return False, (
        f'白色部件分布 左{len(left_parts)}/右{len(right_parts)}/跨中线{len(span_parts)}，'
        f'书脊{len(spine_parts)}件；未构成"书脊+两页对称"或"跨中线自由几何书本"的开书特征'
    )


@rule('D2-25', +1, '左上资源整合图标')
def _(slide, prs):
    """
    细则：位于左上资源整合卡片左上部，距左5.5–7.4cm、距上3.3–5.5cm范围内，
    形状为圆形，直径1.3–1.6cm，填充为橙色，内部为白色打开书本图标。
    检查点（仅按细则约束）：
      1. 圆形位置：left∈[5.5,7.4]cm、top∈[3.3,5.5]cm
      2. 直径：宽高均1.3–1.6cm
      3. 填充为橙色
      4. 圆形内部存在**白色"打开书本"**图标：由 `_is_open_book_icon` 通过
         "左右两页对称"或"跨中线自由几何"两种典型开书几何特征判断，
         不再仅凭"任意 ≥2 个白色部件"通过。
    """
    for circle in slide.shapes:
        l, t, w, h = shape_cm(circle)
        if not in_box(l, t, 5.5, 7.4, 3.3, 5.5):
            continue
        if not is_oval_like(circle):
            continue
        if not (in_range(w, 1.3, 1.6) and in_range(h, 1.3, 1.6)):
            continue
        fill_rgb, _ = get_fill_rgb(circle)
        if not (fill_rgb and is_orange(fill_rgb)):
            continue

        white_parts = []
        for part in slide.shapes:
            if part == circle:
                continue
            if not _shape_inside_circle(part, circle, margin=0.1):
                continue
            if _is_white_editable_icon_part(part):
                white_parts.append(part)

        ok, reason = _is_open_book_icon(circle, white_parts)
        if ok:
            return True, (
                f'{circle.name} 左上资源整合图标 ({l:.2f},{t:.2f}) '
                f'{w:.2f}x{h:.2f}cm 橙色圆形，内部白色打开书本图标：{reason}'
            )
        return False, (
            f'找到橙色圆形图标 {circle.name}，但内部未构成打开书本几何图标：{reason}'
        )

    return False, (
        '未找到符合 left∈[5.5,7.4]cm、top∈[3.3,5.5]cm、'
        '直径1.3–1.6cm、橙色填充的圆形图标'
    )

@rule('D2-26', +1, '左上资源整合标题文本')
def _(slide, prs):
    """
    细则：位于左上资源整合卡片上部，文本为"资源整合"，
    字体为黑体、微软雅黑或相近字体，字号16–18磅，加粗，
    颜色为深橙色，位于图标圆右侧。
    检查点（仅按细则约束）：
      1. 文本精确等于"资源整合"
      2. 位于卡片区域内（left∈[3.0,10.5]cm、top∈[3.0,6.0]cm）
      3. 字体为黑体/微软雅黑或相近中文无衬线字体
      4. 字号16–18磅
      5. 加粗
      6. 颜色为深橙色
      7. 位于图标圆右侧（文本left > 图标圆右边缘 - 0.3cm）
    """
    s = find_shape_with_text(slide, '资源整合')
    if not s:
        return False, '未找到文本精确等于"资源整合"的形状'

    l, t, w, h = shape_cm(s)
    if not in_box(l, t, 3.0, 10.5, 3.0, 6.0):
        return False, f'"资源整合"位置({l:.2f},{t:.2f})不在卡片区域left∈[3.0,10.5]、top∈[3.0,6.0]范围内'

    all_runs = _font_info_all_runs(s)
    if not all_runs:
        fi = font_info(s)
        if fi:
            all_runs = [fi]
    if not all_runs:
        return False, '"资源整合"无法读取字体信息'

    font_ok  = any(is_cjk_sans(r[3]) for r in all_runs)
    sz_ok    = any(r[0] and in_range(r[0], 16, 18) for r in all_runs)
    bold_ok  = any(_is_bold_or_semibold(r[1], r[3]) for r in all_runs)
    color_ok = any(r[2] and (is_deep_orange(r[2]) or is_orange(r[2])) for r in all_runs)

    if not font_ok:
        return False, f'"资源整合"字体{[r[3] for r in all_runs]}非黑体/微软雅黑或相近字体'
    if not sz_ok:
        return False, f'"资源整合"字号{[r[0] for r in all_runs]}不在16–18磅范围'
    if not bold_ok:
        return False, '"资源整合"未加粗'
    if not color_ok:
        return False, f'"资源整合"颜色{[r[2] for r in all_runs]}非深橙色'

    # 位于图标圆右侧：图标圆在 left∈[5.5,7.4]、top∈[3.3,5.5]，直径1.3–1.6cm
    icon_right = None
    for shape in slide.shapes:
        il, it, iw, ih = shape_cm(shape)
        if (is_oval_like(shape) and
                in_box(il, it, 5.5, 7.4, 3.3, 5.5) and
                in_range(iw, 1.3, 1.6) and in_range(ih, 1.3, 1.6)):
            fill_rgb, _ = get_fill_rgb(shape)
            if fill_rgb and is_orange(fill_rgb):
                icon_right = il + iw
                break

    if icon_right is not None and l < icon_right - 0.3:
        return False, (
            f'"资源整合"left={l:.2f}cm 不在图标圆右侧（图标右边缘≈{icon_right:.2f}cm）'
        )

    return True, (
        f'{s.name} "资源整合" 位置({l:.2f},{t:.2f}) '
        f'字号{[r[0] for r in all_runs if r[0]]}pt 加粗 深橙色 图标右侧'
    )

@rule('D2-27', +1, '左上资源整合分隔线与圆点')
def _(slide, prs):
    """
    细则：位于标题文本下方0.3cm–0.7cm处，为一条浅橙色水平单实线，
    长度3.5cm–4.5cm，线宽0.75–1.25磅；中部有一个橙色小圆点。
    检查点（仅按细则约束）：
      1. 以"资源整合"标题文本底边为参考，分隔线中心Y在其下方0.3–0.7cm
      2. 分隔线为浅橙色水平单实线，长度3.5–4.5cm，线宽0.75–1.25磅
      3. 分隔线中部有橙色小圆点（圆形，橙色，中心X/Y与线中心偏差≤0.2cm）
    """
    title_shape = find_shape_with_text(slide, '资源整合')
    if not title_shape:
        return False, '未找到"资源整合"标题文本作为分隔线参考'

    _, title_top, _, title_h = shape_cm(title_shape)
    title_bottom = title_top + title_h
    lines = []
    dots = []

    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        cx, cy = shape_center(s)
        line_rgb = get_line_rgb(s)
        line_width = get_line_width_pt(s)
        if (is_horizontal_line(s, 3.5, 4.5, 0.12) and
                line_rgb and is_light_orange(line_rgb) and
                line_width and in_range(line_width, 0.75, 1.25) and
                not has_dash(s) and
                in_range(cy - title_bottom, 0.3, 0.7)):
            lines.append(s)

        if is_oval_like(s):
            fill_rgb, _ = get_fill_rgb(s)
            if fill_rgb and is_orange(fill_rgb):
                dots.append(s)

    if not lines:
        return False, '未找到位于"资源整合"标题下方0.3–0.7cm、长度3.5–4.5cm、浅橙色、线宽0.75–1.25磅的水平单实线'
    if not dots:
        return False, '未找到橙色小圆点'

    for line in lines:
        lx, ly = line_center(line)
        for dot in dots:
            dx, dy = shape_center(dot)
            if abs(dx - lx) <= 0.2 and abs(dy - ly) <= 0.2:
                ll, lt, lw_cm, lh_cm = shape_cm(line)
                dl, dt, dw, dh = shape_cm(dot)
                return True, (
                    f'{line.name} 位于"资源整合"下方{ly-title_bottom:.2f}cm，'
                    f'长度{lw_cm:.2f}cm 浅橙色水平单实线；'
                    f'{dot.name} 橙色小圆点位于分隔线中部'
                )

    return False, (
        f'找到分隔线{len(lines)}条、橙色小圆点{len(dots)}个，'
        f'但未满足小圆点位于分隔线中部（中心偏差≤0.2cm）'
    )

@rule('D2-28', +1, '左上资源整合三行文本')
def _(slide, prs):
    """
    细则：位于左上资源整合卡片中下部，从上到下出现"图书资源""社区支持""数字素材"，
    字体为黑体、微软雅黑或相近字体，字号13–15磅，颜色为黑色或深灰色，三行均水平居中。
    检查点（仅按细则约束）：
      1. 在卡片区域（left∈[2,10.5]cm、top∈[5.5,7.8]cm）内找到三行文本
      2. 三行按"图书资源"→"社区支持"→"数字素材"自上而下出现
      3. 字体为黑体/微软雅黑或相近中文无衬线字体
      4. 字号13–15磅
      5. 颜色为黑色或深灰色
      6. 三行均水平居中
    """
    lo_l, hi_l, lo_t, hi_t = 2, 10.5, 5.5, 7.8
    required = ('图书资源', '社区支持', '数字素材')

    def line_texts(shape):
        lines = []
        for para in shape.text_frame.paragraphs:
            text = ''.join(r.text for r in para.runs if r.text).strip()
            if text:
                lines.append(text)
        return lines or [shape_text(shape).strip()]

    def all_lines_centered(shape):
        for para in shape.text_frame.paragraphs:
            text = ''.join(r.text for r in para.runs if r.text).strip()
            if not text:
                continue
            align = para.alignment
            ppr = para._p.find('{%s}pPr' % NS['a'])
            xml_center = ppr is not None and ppr.get('algn') == 'ctr'
            if align not in (PP_ALIGN.CENTER, None) and not xml_center:
                return False
        return True

    def style_ok(shape, label):
        all_runs = _font_info_all_runs(shape)
        if not all_runs:
            fi = font_info(shape)
            if fi:
                all_runs = [fi]
        if not all_runs:
            return False, f'"{label}"无法读取字体信息'
        font_ok = any(is_cjk_sans(r[3]) for r in all_runs)
        sz_ok = any(r[0] and in_range(r[0], 13, 15) for r in all_runs)
        color_ok = any(r[2] and _is_dark_or_black(r[2]) for r in all_runs)
        center_ok = all_lines_centered(shape)
        if not font_ok:
            return False, f'"{label}"字体{[r[3] for r in all_runs]}非黑体/微软雅黑或相近字体'
        if not sz_ok:
            return False, f'"{label}"字号{[r[0] for r in all_runs]}不在13–15磅范围'
        if not color_ok:
            return False, f'"{label}"颜色{[r[2] for r in all_runs]}非黑色或深灰色'
        if not center_ok:
            return False, f'"{label}"未水平居中'
        return True, 'ok'

    # Office/WPS 常见做法：三行在同一个文本框内。
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        l, t, w, h = shape_cm(s)
        if not (in_range(l, lo_l, hi_l) and in_range(t, lo_t, hi_t)):
            continue
        lines = line_texts(s)
        if lines == list(required):
            ok, reason = style_ok(s, '、'.join(required))
            if not ok:
                return False, reason
            return True, '三行文本在同一文本框内自上而下出现，字号/字体/颜色/居中均满足细则要求'

    # 兼容三行拆成三个独立文本框的做法。
    found = []
    for text in required:
        s = find_text_in_region(slide, text, lo_l, hi_l, lo_t, hi_t)
        if not s:
            return False, f'在卡片中下部区域未找到文本"{text}"'
        ok, reason = style_ok(s, text)
        if not ok:
            return False, reason
        found.append(s)

    ys = [shape_center(s)[1] for s in found]
    if not all(ys[i] < ys[i + 1] for i in range(len(ys) - 1)):
        return False, f'三行文本未按"图书资源"→"社区支持"→"数字素材"自上而下排列，中心Y={ys}'

    return True, '三行文本分列独立文本框自上而下排列，字号/字体/颜色/居中均满足细则要求'

@rule('D2-29', +3, '左下活动设计卡片外框')
def _(slide, prs):
    """
    细则：位于距左4.3–10.8cm、距上10.5–15.5cm范围内，宽4–5cm，高3.8–4.6cm，
    整体为左侧较宽、右侧略窄的斜边圆角四边形（parallelogram），
    填充为白色，边线为橙色单实线，线宽1.5–2.2磅。
    检查点（仅按细则约束）：
      1. 位置：left∈[4.3,10.8]cm，top∈[10.5,15.5]cm
      2. 尺寸：宽4–5cm，高3.8–4.6cm
      3. 形状为斜边圆角四边形（parallelogram），且**未**水平翻转 → 左侧较宽、右侧略窄
      4. 线条转角为圆角 → 圆角表现
      5. 填充为白色
      6. 边线为橙色单实线，线宽1.5–2.2磅
    """
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        if not in_box(l, t, 4.3, 10.8, 10.5, 15.5):
            continue
        if not in_wh(w, h, 4, 5, 3.8, 4.6):
            continue
        if prst_geom(s) != 'parallelogram':
            continue
        if has_flip_h(s):
            continue
        if not has_round_corners(s):
            continue
        rgb, _ = get_fill_rgb(s)
        if not (rgb is None or is_white_or_near(rgb)):
            continue
        line_rgb = get_line_rgb(s)
        lw = get_line_width_pt(s)
        if not (line_rgb and is_orange(line_rgb)):
            continue
        if not (lw and in_range(lw, 1.5, 2.2)):
            continue
        if has_dash(s):
            continue
        return True, (
            f'{s.name} 左下活动设计卡片外框 ({l:.2f},{t:.2f}) '
            f'{w:.2f}x{h:.2f}cm parallelogram(未翻转→左宽右窄) 圆角 '
            f'白色填充 橙色单实线 线宽{lw:.2f}pt'
        )
    return False, (
        '未找到符合以下全部条件的左下活动设计卡片外框：'
        'left∈[4.3,10.8]cm、top∈[10.5,15.5]cm、宽4–5cm、高3.8–4.6cm、'
        'parallelogram(左宽右窄，未水平翻转)、圆角线条转角、白色填充、'
        '橙色单实线边线、线宽1.5–2.2磅'
    )

def _is_calendar_or_checklist_icon(circle, white_parts):
    """判断圆内白色部件组是否构成"日历"或"任务清单"几何图标。
    典型几何特征（两条路径命中其一）：
      (A) 日历：存在一个"外框"白色矩形/圆角矩形（宽度 ≥ 图标直径 45%，接近方形，
          长宽比 ≤ 2.2），并且外框内至少有 2 条横线或 ≥2 个小格子
          （小矩形/椭圆点/或水平线部件）作为日期格；
      (B) 任务清单：至少 2 条水平线（长宽比 ≥ 3 的横向部件）平行排列（y 差异 ≥ 图标高 15%），
          且每条水平线左侧同一 y 附近存在一个"勾选框/勾/点"部件（矩形/椭圆/
          `checkMark`/长宽比 <2 的小几何）作为清单项。
    该函数不再仅凭"白色部件 ≥ 2"通过；返回 (ok, reason)。"""
    if len(white_parts) < 2:
        return False, f'白色部件仅{len(white_parts)}个'

    _cl, _ct, cw, ch = shape_cm(circle)

    # 收集几何度量
    parts_geo = []  # (part, l, t, w, h, cx, cy, ratio, geom)
    for p in white_parts:
        pl, pt, pw, ph = shape_cm(p)
        if pw <= 0 or ph <= 0:
            pw = max(pw, 0.001); ph = max(ph, 0.001)
        ratio = max(pw, ph) / max(min(pw, ph), 0.001)
        parts_geo.append((p, pl, pt, pw, ph, pl + pw / 2, pt + ph / 2, ratio, prst_geom(p)))

    # ── (A) 日历：外框 + 日期格/横线 ─────────────────────────────
    frame_candidates = []
    for g in parts_geo:
        _p, _pl, _pt, pw, _ph, _pcx, _pcy, ratio, geom = g
        if geom in ('rect', 'roundRect') and pw >= cw * 0.45 and ratio <= 2.2:
            frame_candidates.append(g)

    if frame_candidates:
        fr = max(frame_candidates, key=lambda g: g[3])
        fp, fl, ft, fw, fh, _fcx, _fcy, _fr_ratio, _fr_geom = fr
        cells = 0
        hlines = 0
        for g in parts_geo:
            if g[0] is fp:
                continue
            _p2, _l2, _t2, _w2, _h2, _cx2, _cy2, _r2, _geom2 = g
            if not (fl - 0.05 <= _cx2 <= fl + fw + 0.05):
                continue
            if not (ft - 0.05 <= _cy2 <= ft + fh + 0.05):
                continue
            if _r2 >= 3.0 and _w2 >= fw * 0.35:
                hlines += 1
            elif _r2 < 2.5 and _w2 <= fw * 0.5 and _h2 <= fh * 0.5:
                cells += 1
        if hlines >= 2 or cells >= 2:
            return True, (
                f'日历图标：外框 {fw:.2f}x{fh:.2f}cm + '
                f'{hlines}条横线/{cells}个日期格'
            )

    # ── (B) 任务清单：≥2 条平行横线 + 每条左侧勾选框 ─────────────
    hlines2 = []
    for g in parts_geo:
        _p, _pl, _pt, pw, ph, _pcx, pcy, ratio, _geom = g
        if ratio >= 3.0 and pw >= cw * 0.35 and ph <= ch * 0.2:
            hlines2.append((g[0], g[1], g[2], pw, ph, pcy))
    if len(hlines2) >= 2:
        hlines_sorted = sorted(hlines2, key=lambda t: t[5])
        parallel_pairs = [(a, b) for i, a in enumerate(hlines_sorted)
                          for b in hlines_sorted[i + 1:]
                          if b[5] - a[5] >= ch * 0.15]
        if parallel_pairs:
            def _has_checkbox_left_of(line):
                _lp, ll, _lt, _lw, _lh, lcy = line
                for g in parts_geo:
                    p2, _pl2, _pt2, pw2, ph2, pcx2, pcy2, r2, geom2 = g
                    if p2 is line[0]:
                        continue
                    if geom2 == 'checkMark':
                        if pcx2 < ll and abs(pcy2 - lcy) <= ch * 0.12:
                            return True
                        continue
                    if r2 < 2.0 and pw2 <= cw * 0.28 and ph2 <= ch * 0.28:
                        if pcx2 < ll + 0.02 and abs(pcy2 - lcy) <= ch * 0.15:
                            return True
                return False
            ok_lines = [ln for ln in hlines2 if _has_checkbox_left_of(ln)]
            if len(ok_lines) >= 2:
                return True, (
                    f'任务清单图标：{len(hlines2)}条平行横线，'
                    f'{len(ok_lines)}条左侧含勾选框/勾/点'
                )

    return False, (
        f'白色部件{len(white_parts)}个未构成日历(外框+日期格/横线)'
        f'或任务清单(≥2条平行横线+左侧勾选框)几何图标'
    )


@rule('D2-30', +1, '左下活动设计图标')
def _(slide, prs):
    """
    细则：位于左下活动设计卡片左上部，距左5.2–7.7cm、距上11.0–12.8cm范围内，
    形状为圆形，直径1.3–1.6cm，填充为橙色，内部为白色日历或任务清单图标。
    检查点（仅按细则约束）：
      1. 圆形位置：left∈[5.2,7.7]cm、top∈[11.0,12.8]cm
      2. 直径：宽高均1.3–1.6cm
      3. 填充为橙色
      4. 圆形内部存在**白色"日历或任务清单"**图标：由 `_is_calendar_or_checklist_icon`
         按"外框+日期格/横线"或"≥2条平行横线+左侧勾选框"两种典型结构特征判断，
         不再仅凭"任意 ≥2 个白色部件"通过。
    """
    for circle in slide.shapes:
        l, t, w, h = shape_cm(circle)
        if not in_box(l, t, 5.2, 7.7, 11.0, 12.8):
            continue
        if not is_oval_like(circle):
            continue
        if not (in_range(w, 1.3, 1.6) and in_range(h, 1.3, 1.6)):
            continue
        fill_rgb, _ = get_fill_rgb(circle)
        if not (fill_rgb and is_orange(fill_rgb)):
            continue

        white_parts = []
        for part in slide.shapes:
            if part == circle:
                continue
            if not _shape_inside_circle(part, circle, margin=0.1):
                continue
            if _is_white_editable_icon_part(part):
                white_parts.append(part)

        ok, reason = _is_calendar_or_checklist_icon(circle, white_parts)
        if ok:
            return True, (
                f'{circle.name} 左下活动设计图标 ({l:.2f},{t:.2f}) '
                f'{w:.2f}x{h:.2f}cm 橙色圆形，内部白色日历/任务清单图标：{reason}'
            )
        return False, (
            f'找到橙色圆形图标 {circle.name}，但内部未构成日历/任务清单几何图标：{reason}'
        )

    return False, (
        '未找到符合 left∈[5.2,7.7]cm、top∈[11.0,12.8]cm、'
        '直径1.3–1.6cm、橙色填充的圆形图标'
    )

@rule('D2-31', +1, '左下活动设计标题文本')
def _(slide, prs):
    """
    细则：位于左下活动设计卡片上部，文本为"活动设计"，
    字体为黑体、微软雅黑或相近字体，字号16–18磅，加粗，
    颜色为深橙色，位于图标圆右侧。
    检查点（仅按细则约束）：
      1. 文本精确等于"活动设计"
      2. 位于左下活动设计卡片上部（left∈[3.0,10.8]cm、top∈[11.0,13.8]cm）
      3. 字体为黑体/微软雅黑或相近中文无衬线字体
      4. 字号16–18磅
      5. 加粗
      6. 颜色为深橙色
      7. 位于图标圆右侧（文本left > 图标圆右边缘 - 0.3cm）
    """
    s = find_shape_with_text(slide, '活动设计')
    if not s:
        return False, '未找到文本精确等于"活动设计"的形状'

    l, t, w, h = shape_cm(s)
    if not in_box(l, t, 3.0, 10.8, 11.0, 13.8):
        return False, f'"活动设计"位置({l:.2f},{t:.2f})不在卡片上部left∈[3.0,10.8]、top∈[11.0,13.8]范围内'

    all_runs = _font_info_all_runs(s)
    if not all_runs:
        fi = font_info(s)
        if fi:
            all_runs = [fi]
    if not all_runs:
        return False, '"活动设计"无法读取字体信息'

    font_ok  = any(is_cjk_sans(r[3]) for r in all_runs)
    sz_ok    = any(r[0] and in_range(r[0], 16, 18) for r in all_runs)
    bold_ok  = any(_is_bold_or_semibold(r[1], r[3]) for r in all_runs)
    color_ok = any(r[2] and (is_deep_orange(r[2]) or is_orange(r[2])) for r in all_runs)

    if not font_ok:
        return False, f'"活动设计"字体{[r[3] for r in all_runs]}非黑体/微软雅黑或相近字体'
    if not sz_ok:
        return False, f'"活动设计"字号{[r[0] for r in all_runs]}不在16–18磅范围'
    if not bold_ok:
        return False, '"活动设计"未加粗'
    if not color_ok:
        return False, f'"活动设计"颜色{[r[2] for r in all_runs]}非深橙色'

    icon_right = None
    for shape in slide.shapes:
        il, it, iw, ih = shape_cm(shape)
        if (is_oval_like(shape) and
                in_box(il, it, 5.2, 7.7, 11.0, 12.8) and
                in_range(iw, 1.3, 1.6) and in_range(ih, 1.3, 1.6)):
            fill_rgb, _ = get_fill_rgb(shape)
            if fill_rgb and is_orange(fill_rgb):
                icon_right = il + iw
                break

    if icon_right is not None and l < icon_right - 0.3:
        return False, (
            f'"活动设计"left={l:.2f}cm 不在图标圆右侧（图标右边缘≈{icon_right:.2f}cm）'
        )

    return True, (
        f'{s.name} "活动设计" 位置({l:.2f},{t:.2f}) '
        f'字号{[r[0] for r in all_runs if r[0]]}pt 加粗 深橙色 图标右侧'
    )

@rule('D2-32', +1, '左下活动设计分隔线与圆点')
def _(slide, prs):
    """
    细则：位于标题文本下方0.3cm–0.7cm处，为一条浅橙色水平单实线，
    长度3.5cm–4.5cm，线宽0.75–1.25磅；中部有一个橙色小圆点。
    检查点（仅按细则约束）：
      1. 以"活动设计"标题文本底边为参考，分隔线中心Y在其下方0.3–0.7cm
      2. 分隔线为浅橙色水平单实线，长度3.5–4.5cm，线宽0.75–1.25磅
      3. 分隔线中部有橙色小圆点（圆形，橙色，中心X/Y与线中心偏差≤0.2cm）
    """
    title_shape = find_shape_with_text(slide, '活动设计')
    if not title_shape:
        return False, '未找到"活动设计"标题文本作为分隔线参考'

    _, title_top, _, title_h = shape_cm(title_shape)
    title_bottom = title_top + title_h
    lines = []
    dots = []

    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        line_rgb = get_line_rgb(s)
        line_width = get_line_width_pt(s)
        _, cy = shape_center(s)
        if (is_horizontal_line(s, 3.5, 4.5, 0.12) and
                line_rgb and is_light_orange(line_rgb) and
                line_width and in_range(line_width, 0.75, 1.25) and
                not has_dash(s) and
                in_range(cy - title_bottom, 0.3, 0.7)):
            lines.append(s)

        if is_oval_like(s):
            fill_rgb, _ = get_fill_rgb(s)
            if fill_rgb and is_orange(fill_rgb):
                dots.append(s)

    if not lines:
        return False, '未找到位于"活动设计"标题下方0.3–0.7cm、长度3.5–4.5cm、浅橙色、线宽0.75–1.25磅的水平单实线'
    if not dots:
        return False, '未找到橙色小圆点'

    for line in lines:
        lx, ly = line_center(line)
        for dot in dots:
            dx, dy = shape_center(dot)
            if abs(dx - lx) <= 0.2 and abs(dy - ly) <= 0.2:
                ll, lt, lw_cm, lh_cm = shape_cm(line)
                return True, (
                    f'{line.name} 位于"活动设计"下方{ly-title_bottom:.2f}cm，'
                    f'长度{lw_cm:.2f}cm 浅橙色水平单实线；'
                    f'{dot.name} 橙色小圆点位于分隔线中部'
                )

    return False, (
        f'找到分隔线{len(lines)}条、橙色小圆点{len(dots)}个，'
        f'但未满足小圆点位于分隔线中部（中心偏差≤0.2cm）'
    )

@rule('D2-33', +1, '左下活动设计三行文本')
def _(slide, prs):
    """
    细则：位于左下活动设计卡片中下部，从上到下出现"晨读打卡""主题书展""读书分享"，
    字体为黑体、微软雅黑或相近字体，字号15–21磅，颜色为黑色或深灰色，三行均水平居中。
    检查点（仅按细则约束）：
      1. 在卡片中下部区域（left∈[2,10.8]cm、top∈[13.0,15.5]cm）内找到三行文本
      2. 三行按"晨读打卡"→"主题书展"→"读书分享"自上而下排列（中心Y递增）
      3. 每行字体为黑体/微软雅黑或相近中文无衬线字体
      4. 每行字号15–21磅
      5. 每行颜色为黑色或深灰色
      6. 每行水平居中
    """
    lo_l, hi_l, lo_t, hi_t = 2, 10.8, 13.0, 15.5
    required = ('晨读打卡', '主题书展', '读书分享')
    found = []

    for text in required:
        s = find_text_in_region(slide, text, lo_l, hi_l, lo_t, hi_t)
        if not s:
            return False, f'在左下活动设计卡片中下部区域未找到文本"{text}"'

        all_runs = _font_info_all_runs(s)
        if not all_runs:
            fi = font_info(s)
            if fi:
                all_runs = [fi]
        if not all_runs:
            return False, f'"{text}"无法读取字体信息'

        font_ok  = any(is_cjk_sans(r[3]) for r in all_runs)
        sz_ok    = any(r[0] and in_range(r[0], 15, 21) for r in all_runs)
        color_ok = any(r[2] and _is_dark_or_black(r[2]) for r in all_runs)
        center_ok = text_is_centered(s)

        if not font_ok:
            return False, f'"{text}"字体{[r[3] for r in all_runs]}非黑体/微软雅黑或相近字体'
        if not sz_ok:
            return False, f'"{text}"字号{[r[0] for r in all_runs]}不在15–21磅范围'
        if not color_ok:
            return False, f'"{text}"颜色{[r[2] for r in all_runs]}非黑色或深灰色'
        if not center_ok:
            return False, f'"{text}"未水平居中'

        found.append(s)

    ys = [shape_center(s)[1] for s in found]
    if not all(ys[i] < ys[i+1] for i in range(len(ys)-1)):
        return False, f'三行文本未按"晨读打卡"→"主题书展"→"读书分享"自上而下排列，中心Y={ys}'

    return True, (
        f'"晨读打卡""主题书展""读书分享"均在卡片中下部，'
        f'自上而下排列，字号/字体/颜色/居中均满足细则要求'
    )

@rule('D2-34', +3, '右上教学支持卡片外框')
def _(slide, prs):
    """
    细则：位于距左23–29cm、距上3.1–7.8cm范围内，宽4–5cm，高3.8–4.6cm，
    整体为右侧较宽、左侧略窄的斜边圆角四边形（parallelogram，水平翻转），
    填充为白色，边线为橙色单实线，线宽1.5–2.2磅。
    检查点（仅按细则约束）：
      1. 位置：left∈[23,29]cm，top∈[3.1,7.8]cm
      2. 尺寸：宽4–5cm，高3.8–4.6cm
      3. 形状为斜边圆角四边形（parallelogram），且**水平翻转** → 右侧较宽、左侧略窄
      4. 线条转角为圆角 → 圆角表现
      5. 填充为白色
      6. 边线为橙色单实线，线宽1.5–2.2磅
    """
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        if not in_box(l, t, 23, 29, 3.1, 7.8):
            continue
        if not in_wh(w, h, 4, 5, 3.8, 4.6):
            continue
        if prst_geom(s) != 'parallelogram':
            continue
        # 方向：右宽左窄 ⇒ parallelogram 水平翻转
        if not has_flip_h(s):
            continue
        if not has_round_corners(s):
            continue
        rgb, _ = get_fill_rgb(s)
        if not (rgb is None or is_white_or_near(rgb)):
            continue
        line_rgb = get_line_rgb(s)
        lw = get_line_width_pt(s)
        if not (line_rgb and is_orange(line_rgb)):
            continue
        if not (lw and in_range(lw, 1.5, 2.2)):
            continue
        if has_dash(s):
            continue
        return True, (
            f'{s.name} 右上教学支持卡片外框 ({l:.2f},{t:.2f}) '
            f'{w:.2f}x{h:.2f}cm parallelogram(水平翻转→右宽左窄) 圆角 '
            f'白色填充 橙色单实线 线宽{lw:.2f}pt'
        )
    return False, (
        '未找到符合以下全部条件的右上教学支持卡片外框：'
        'left∈[23,29]cm、top∈[3.1,7.8]cm、宽4–5cm、高3.8–4.6cm、'
        'parallelogram(右宽左窄，水平翻转)、圆角线条转角、白色填充、'
        '橙色单实线边线、线宽1.5–2.2磅'
    )

def _is_graduation_cap_icon(circle, white_parts):
    """判断圆内白色部件组是否构成"学位帽（mortarboard）"几何图标。
    学位帽典型结构（择一命中即可）：
      (A) 帽板+帽体两段式：
          - 帽板（顶部）：一件"扁而宽"的四边形，长宽比 ≥ 2、宽 ≥ 图标直径 55%，
            几何为 parallelogram/diamond/rhombus 类或近似菱形/斜四边形；
          - 帽体（下方）：一件"梯形/曲边"部件在帽板下方，中心 y > 帽板中心 y；
            几何可为 trapezoid/`homePlate`/`plaque`/`arc`/`curvedX`/`custGeom` 或
            近似矩形/宽扁形状（宽 ≥ 图标直径 40%）；
      (B) 帽板 + 流苏（tassel）：帽板同 (A)，且从帽板一侧向下悬挂一件"细长竖直"
          部件（长宽比 ≥ 3、height ≥ 图标高 25%），代表流苏/穗。
      (C) 单件整体形状：整顶学位帽由一件 parallelogram/diamond/rhombus
          或 custGeom/freeform 承载，无独立帽体/流苏部件——
          该形状须占据图标显著面积（宽 ≥ 图标直径 55%、高 ≥ 图标高 35%），
          具备"扁而宽"的帽形（长宽比 ≥ 1.6），且位于图标中央附近。
      (D) 纯线条骨架：整顶学位帽由多条白色线条（`line`/`straightConnector1`/
          `bentConnector*`/`curvedConnector*` 或极瘦长自定义形状）勾勒——
          须 ≥ 3 条白线，其中至少 1 条水平线（宽 ≥ 图标直径 55%）作为帽板顶边，
          下方另有 ≥ 2 条线（斜边/竖线）组成帽体或流苏骨架；线条整体外包围盒
          宽 ≥ 图标直径 55%、高 ≥ 图标高 30%，且居于图标中央。
    该函数不再仅凭"白色部件 ≥ 2"通过；返回 (ok, reason)。"""
    if len(white_parts) < 1:
        return False, '未在圆内找到白色部件'

    _cl, _ct, cw, ch = shape_cm(circle)
    ccx = _cl + cw / 2.0
    ccy = _ct + ch / 2.0

    parts_geo = []  # (part, l, t, w, h, cx, cy, ratio, geom)
    for p in white_parts:
        pl, pt, pw, ph = shape_cm(p)
        if pw <= 0 or ph <= 0:
            pw = max(pw, 0.001); ph = max(ph, 0.001)
        ratio = max(pw, ph) / max(min(pw, ph), 0.001)
        parts_geo.append((p, pl, pt, pw, ph, pl + pw / 2, pt + ph / 2, ratio, prst_geom(p)))

    # ── (C) 单件整体学位帽：parallelogram/diamond/rhombus 或 custGeom/freeform ─
    SINGLE_CAP_GEOMS = {'parallelogram', 'diamond', 'rhombus'}
    for g in parts_geo:
        p, _pl, _pt, pw, ph, pcx, pcy, ratio, geom = g
        is_custom = (geom in ('freeform', '')
                     and p._element.find('.//{%s}custGeom' % NS['a']) is not None) \
                    or geom == 'freeform'
        if not (geom in SINGLE_CAP_GEOMS or is_custom):
            continue
        # 占据图标显著面积、扁而宽、位于图标中央
        if pw < cw * 0.55 or ph < ch * 0.35:
            continue
        if ratio < 1.6:
            continue
        if abs(pcx - ccx) > cw * 0.25 or abs(pcy - ccy) > ch * 0.35:
            continue
        return True, (
            f'学位帽(单件整体 {geom or "custGeom"})：{pw:.2f}x{ph:.2f}cm '
            f'长宽比{ratio:.2f} 位于图标中央'
        )

    if len(white_parts) < 2:
        return False, (
            f'白色部件仅{len(white_parts)}个：既不构成(C)单件整体学位帽，'
            f'也不足以组成(A)帽板+帽体或(B)帽板+流苏结构'
        )

    # ── 帽板候选（扁而宽的四边形，靠近图标上半部）─────────────────
    CAP_PLATE_GEOMS = {
        'parallelogram', 'diamond', 'rhombus', 'trapezoid',
        'rect', 'roundRect', 'flowChartConnector',
    }
    plate_candidates = []
    for g in parts_geo:
        _p, _pl, _pt, pw, ph, _pcx, pcy, ratio, geom = g
        # 扁而宽 & 尺寸达标 & 靠近上半部
        if pw >= cw * 0.55 and ph <= ch * 0.6 and ratio >= 2.0 and pcy <= ccy + ch * 0.1:
            if geom in CAP_PLATE_GEOMS or geom == '' or geom.startswith('curved'):
                plate_candidates.append(g)

    # ── (D) 纯线条骨架学位帽（无闭合帽板/帽体也可命中）──────────────
    # 由若干条白色线条勾勒出的学位帽：至少一条水平线充当帽板顶边，
    # 下方另有 ≥2 条线（斜边/竖线）组成帽体或流苏骨架；线条整体
    # 外包围盒宽 ≥ 图标直径 55%、高 ≥ 图标高 30%、居于图标中央。
    # 说明：因无合格样例，阈值按常识保守；后续可依样例微调。
    LINE_GEOMS = {
        'line', 'straightConnector1',
        'bentConnector2', 'bentConnector3', 'bentConnector4', 'bentConnector5',
        'curvedConnector2', 'curvedConnector3', 'curvedConnector4', 'curvedConnector5',
    }

    def _is_line_like(part, pw_, ph_, geom_):
        if geom_ in LINE_GEOMS:
            return True
        thin = min(pw_, ph_) <= max(0.06, 0.04 * max(cw, ch))
        long_ = max(pw_, ph_) >= 0.12 * max(cw, ch)
        return thin and long_

    line_parts = []  # (part, l, t, w, h, cx, cy, ratio, geom)
    for g in parts_geo:
        _p, _pl, _pt, pw_, ph_, _pcx, _pcy, _r, geom_ = g
        if _is_line_like(_p, pw_, ph_, geom_):
            line_parts.append(g)

    if len(line_parts) >= 3:
        horiz = None
        for g in line_parts:
            _p, _pl, _pt, pw_, ph_, _pcx, _pcy, _r, _geom = g
            if pw_ >= cw * 0.55 and ph_ <= max(0.08, ch * 0.08) and pw_ > ph_ * 4:
                if horiz is None or pw_ > horiz[3]:
                    horiz = g
        if horiz is not None:
            _hp, _hl, _ht, _hw, _hh, _hcx, hcy, _hr, _hg = horiz
            below = 0
            for g in line_parts:
                _p, _pl, _pt, pw_, ph_, _pcx, _pcy, _r, _geom = g
                if _p is _hp:
                    continue
                if _pcy > hcy + max(0.05, ch * 0.05):
                    below += 1
            if below >= 2:
                ux0 = min(g[1] for g in line_parts)
                uy0 = min(g[2] for g in line_parts)
                ux1 = max(g[1] + g[3] for g in line_parts)
                uy1 = max(g[2] + g[4] for g in line_parts)
                uw, uh = ux1 - ux0, uy1 - uy0
                ucx = (ux0 + ux1) / 2.0
                ucy = (uy0 + uy1) / 2.0
                if (uw >= cw * 0.55 and uh >= ch * 0.30
                        and abs(ucx - ccx) <= cw * 0.25
                        and abs(ucy - ccy) <= ch * 0.35):
                    return True, (
                        f'学位帽(线条骨架)：{len(line_parts)}条白色线条，'
                        f'顶边水平线宽{_hw:.2f}cm，下方{below}条支撑线，'
                        f'整体{uw:.2f}x{uh:.2f}cm 居中'
                    )

    if not plate_candidates:
        return False, (
            f'白色部件{len(white_parts)}个未找到学位帽帽板：'
            f'需一件宽≥{cw*0.55:.2f}cm、长宽比≥2、位于上半部的扁四边形；'
            f'亦不满足纯线条骨架条件'
        )

    plate = max(plate_candidates, key=lambda g: g[3])  # 取最宽者
    pp, pl, pt, pw, ph, pcx, pcy, _p_ratio, _p_geom = plate

    # ── (A) 帽体：帽板下方的宽扁部件 ──────────────────────────────
    CAP_BODY_GEOMS = {
        'trapezoid', 'homePlate', 'plaque', 'rect', 'roundRect',
        'arc',
    }
    for g in parts_geo:
        p2, _l2, _t2, w2, h2, _cx2, cy2, r2, geom2 = g
        if p2 is pp:
            continue
        # 位于帽板正下方（中心 y 在帽板下方且距离不太远）
        if not (pcy < cy2 <= pcy + ch * 0.9):
            continue
        # 水平中心与帽板中心大致对齐（容差为图标宽 25%）
        if abs(_cx2 - pcx) > cw * 0.25:
            continue
        # 宽扁：宽 ≥ 图标 40%
        if w2 >= cw * 0.40 and (
            geom2 in CAP_BODY_GEOMS or geom2.startswith('curved')
            or geom2 == '' or p2._element.find('.//{%s}custGeom' % NS['a']) is not None
        ):
            return True, (
                f'学位帽(帽板+帽体)：帽板 {pw:.2f}x{ph:.2f}cm({_p_geom or "custGeom"}) '
                f'+ 帽体 {w2:.2f}x{h2:.2f}cm({geom2 or "custGeom"})'
            )

    # ── (B) 帽板 + 流苏 ────────────────────────────────────────
    for g in parts_geo:
        p2, _l2, t2, w2, h2, cx2, cy2, r2, _geom2 = g
        if p2 is pp:
            continue
        # 细长竖直：长宽比 ≥ 3 且 height ≥ 图标高 25%
        if r2 >= 3.0 and h2 >= ch * 0.25 and h2 > w2:
            # 位于帽板一侧且向下延伸：cy 在帽板下方
            if cy2 > pcy and abs(cx2 - pcx) <= pw * 0.6:
                return True, (
                    f'学位帽(帽板+流苏)：帽板 {pw:.2f}x{ph:.2f}cm '
                    f'+ 流苏 {w2:.2f}x{h2:.2f}cm(细长竖直)'
                )

    return False, (
        f'白色部件{len(white_parts)}个：找到帽板 {pw:.2f}x{ph:.2f}cm，'
        f'但下方未找到帽体(宽扁梯形/曲边)或流苏(细长竖直)，'
        f'亦不满足纯线条骨架条件(≥3条白线、含≥1条水平顶边、'
        f'≥2条下方支撑、整体宽高居中)'
    )


@rule('D2-35', +1, '右上教学支持图标')
def _(slide, prs):
    """
    细则：位于右上教学支持卡片左上部，距左24.3–26cm、距上3.3–5.5cm范围内，
    形状为圆形，直径1.3–1.6cm，填充为橙色，内部为白色学位帽图标。
    检查点（仅按细则约束）：
      1. 圆形位置：left∈[24.3,26.0]cm、top∈[3.3,5.5]cm
      2. 直径：宽高均1.3–1.6cm
      3. 填充为橙色
      4. 圆形内部存在**白色"学位帽"**图标：由 `_is_graduation_cap_icon`
         按"帽板+帽体"或"帽板+流苏"两种典型结构特征判断，
         不再仅凭"任意 ≥2 个白色部件"通过。
    """
    for circle in slide.shapes:
        l, t, w, h = shape_cm(circle)
        if not in_box(l, t, 24.3, 26.0, 3.3, 5.5):
            continue
        if not is_oval_like(circle):
            continue
        if not (in_range(w, 1.3, 1.6) and in_range(h, 1.3, 1.6)):
            continue
        fill_rgb, _ = get_fill_rgb(circle)
        if not (fill_rgb and is_orange(fill_rgb)):
            continue

        white_parts = []
        for part in slide.shapes:
            if part == circle:
                continue
            if not _shape_inside_circle(part, circle, margin=0.1):
                continue
            if _is_white_editable_icon_part(part):
                white_parts.append(part)

        ok, reason = _is_graduation_cap_icon(circle, white_parts)
        if ok:
            return True, (
                f'{circle.name} 右上教学支持图标 ({l:.2f},{t:.2f}) '
                f'{w:.2f}x{h:.2f}cm 橙色圆形，内部白色学位帽图标：{reason}'
            )
        return False, (
            f'找到橙色圆形图标 {circle.name}，但内部未构成学位帽几何图标：{reason}'
        )

    return False, (
        '未找到符合 left∈[24.3,26.0]cm、top∈[3.3,5.5]cm、'
        '直径1.3–1.6cm、橙色填充的圆形图标'
    )

@rule('D2-36', +1, '右上教学支持标题文本')
def _(slide, prs):
    """
    细则：位于右上教学支持卡片上部，文本为"教学支持"，
    字体为黑体、微软雅黑或相近字体，字号16–18磅，加粗，
    颜色为深橙色，位于图标圆右侧。
    检查点（仅按细则约束）：
      1. 文本精确等于"教学支持"
      2. 位于卡片上部（left∈[26.0,33.0]cm、top∈[3.0,6.0]cm）
      3. 字体为黑体/微软雅黑或相近中文无衬线字体
      4. 字号16–18磅
      5. 加粗
      6. 颜色为深橙色
      7. 位于图标圆右侧（文本left > 图标圆右边缘 - 0.3cm）
    """
    s = find_shape_with_text(slide, '教学支持')
    if not s:
        return False, '未找到文本精确等于"教学支持"的形状'

    l, t, w, h = shape_cm(s)
    if not in_box(l, t, 26.0, 33.0, 3.0, 6.0):
        return False, f'"教学支持"位置({l:.2f},{t:.2f})不在卡片上部left∈[26.0,33.0]、top∈[3.0,6.0]范围内'

    all_runs = _font_info_all_runs(s)
    if not all_runs:
        fi = font_info(s)
        if fi:
            all_runs = [fi]
    if not all_runs:
        return False, '"教学支持"无法读取字体信息'

    font_ok  = any(is_cjk_sans(r[3]) for r in all_runs)
    sz_ok    = any(r[0] and in_range(r[0], 16, 18) for r in all_runs)
    bold_ok  = any(_is_bold_or_semibold(r[1], r[3]) for r in all_runs)
    color_ok = any(r[2] and (is_deep_orange(r[2]) or is_orange(r[2])) for r in all_runs)

    if not font_ok:
        return False, f'"教学支持"字体{[r[3] for r in all_runs]}非黑体/微软雅黑或相近字体'
    if not sz_ok:
        return False, f'"教学支持"字号{[r[0] for r in all_runs]}不在16–18磅范围'
    if not bold_ok:
        return False, '"教学支持"未加粗'
    if not color_ok:
        return False, f'"教学支持"颜色{[r[2] for r in all_runs]}非深橙色'

    # 位于图标圆右侧：图标圆在 left∈[24.3,26.0]、top∈[3.3,5.5]，直径1.3–1.6cm
    icon_right = None
    for shape in slide.shapes:
        il, it, iw, ih = shape_cm(shape)
        if (is_oval_like(shape) and
                in_box(il, it, 24.3, 26.0, 3.3, 5.5) and
                in_range(iw, 1.3, 1.6) and in_range(ih, 1.3, 1.6)):
            fill_rgb, _ = get_fill_rgb(shape)
            if fill_rgb and is_orange(fill_rgb):
                icon_right = il + iw
                break

    if icon_right is not None and l < icon_right - 0.3:
        return False, (
            f'"教学支持"left={l:.2f}cm 不在图标圆右侧（图标右边缘≈{icon_right:.2f}cm）'
        )

    return True, (
        f'{s.name} "教学支持" 位置({l:.2f},{t:.2f}) '
        f'字号{[r[0] for r in all_runs if r[0]]}pt 加粗 深橙色 图标右侧'
    )

@rule('D2-37', +1, '右上教学支持分隔线与圆点')
def _(slide, prs):
    """
    细则：位于标题文本下方0.3cm–0.7cm处，为一条浅橙色水平单实线，
    长度3.5cm–4.5cm，线宽0.75–1.25磅；中部有一个橙色小圆点。
    检查点（仅按细则约束）：
      1. 以"教学支持"标题文本底边为参考，分隔线中心Y在其下方0.3–0.7cm
      2. 分隔线为浅橙色水平单实线，长度3.5–4.5cm，线宽0.75–1.25磅
      3. 分隔线中部有橙色小圆点（圆形，橙色，中心X/Y与线中心偏差≤0.2cm）
    """
    title_shape = find_shape_with_text(slide, '教学支持')
    if not title_shape:
        return False, '未找到"教学支持"标题文本作为分隔线参考'

    _, title_top, _, title_h = shape_cm(title_shape)
    title_bottom = title_top + title_h
    lines = []
    dots = []

    for s in slide.shapes:
        line_rgb = get_line_rgb(s)
        line_width = get_line_width_pt(s)
        _, cy = shape_center(s)
        if (is_horizontal_line(s, 3.5, 4.5, 0.12) and
                line_rgb and is_light_orange(line_rgb) and
                line_width and in_range(line_width, 0.75, 1.25) and
                not has_dash(s) and
                in_range(cy - title_bottom, 0.3, 0.7)):
            lines.append(s)

        if is_oval_like(s):
            fill_rgb, _ = get_fill_rgb(s)
            if fill_rgb and is_orange(fill_rgb):
                dots.append(s)

    if not lines:
        return False, '未找到位于"教学支持"标题下方0.3–0.7cm、长度3.5–4.5cm、浅橙色、线宽0.75–1.25磅的水平单实线'
    if not dots:
        return False, '未找到橙色小圆点'

    for line in lines:
        lx, ly = line_center(line)
        for dot in dots:
            dx, dy = shape_center(dot)
            if abs(dx - lx) <= 0.2 and abs(dy - ly) <= 0.2:
                ll, lt, lw_cm, lh_cm = shape_cm(line)
                return True, (
                    f'{line.name} 位于"教学支持"下方{ly-title_bottom:.2f}cm，'
                    f'长度{lw_cm:.2f}cm 浅橙色水平单实线；'
                    f'{dot.name} 橙色小圆点位于分隔线中部'
                )

    return False, (
        f'找到分隔线{len(lines)}条、橙色小圆点{len(dots)}个，'
        f'但未满足小圆点位于分隔线中部（中心偏差≤0.2cm）'
    )

@rule('D2-38', +1, '右上教学支持三行文本')
def _(slide, prs):
    """
    细则：位于右上教学支持卡片中下部，从上到下出现"分层任务""情境引导""跨学科融合"，
    字体为黑体、微软雅黑或相近字体，字号13–15磅，颜色为黑色或深灰色，三行均水平居中。
    """
    lo_l, hi_l, lo_t, hi_t = 24, 32, 5.5, 7.8
    required = ('分层任务', '情境引导', '跨学科融合')

    def line_texts(shape):
        lines = []
        for para in shape.text_frame.paragraphs:
            text = ''.join(r.text for r in para.runs if r.text).strip()
            if text:
                lines.append(text)
        return lines or [shape_text(shape).strip()]

    def all_lines_centered(shape):
        for para in shape.text_frame.paragraphs:
            text = ''.join(r.text for r in para.runs if r.text).strip()
            if not text:
                continue
            align = para.alignment
            ppr = para._p.find('{%s}pPr' % NS['a'])
            xml_center = ppr is not None and ppr.get('algn') == 'ctr'
            if align not in (PP_ALIGN.CENTER, None) and not xml_center:
                return False
        return True

    def style_ok(shape, label):
        all_runs = _font_info_all_runs(shape)
        if not all_runs:
            fi = font_info(shape)
            if fi:
                all_runs = [fi]
        if not all_runs:
            return False, f'"{label}"无法读取字体信息'
        font_ok = any(is_cjk_sans(r[3]) for r in all_runs)
        sz_ok = any(r[0] and in_range(r[0], 13, 15) for r in all_runs)
        color_ok = any(r[2] and _is_dark_or_black(r[2]) for r in all_runs)
        center_ok = all_lines_centered(shape)
        if not font_ok:
            return False, f'"{label}"字体{[r[3] for r in all_runs]}非黑体/微软雅黑或相近字体'
        if not sz_ok:
            return False, f'"{label}"字号{[r[0] for r in all_runs]}不在13–15磅范围'
        if not color_ok:
            return False, f'"{label}"颜色{[r[2] for r in all_runs]}非黑色或深灰色'
        if not center_ok:
            return False, f'"{label}"未水平居中'
        return True, 'ok'

    # 同一文本框三行
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        l, t, w, h = shape_cm(s)
        if not (in_range(l, lo_l, hi_l) and in_range(t, lo_t, hi_t)):
            continue
        if line_texts(s) == list(required):
            ok, reason = style_ok(s, '、'.join(required))
            if not ok:
                return False, reason
            return True, '三行文本在同一文本框内自上而下出现，字号/字体/颜色/居中均满足细则要求'

    # 三个独立文本框
    found = []
    for text in required:
        s = find_text_in_region(slide, text, lo_l, hi_l, lo_t, hi_t)
        if not s:
            return False, f'在右上教学支持卡片中下部区域未找到文本"{text}"'
        ok, reason = style_ok(s, text)
        if not ok:
            return False, reason
        found.append(s)

    ys = [shape_center(s)[1] for s in found]
    if not all(ys[i] < ys[i + 1] for i in range(len(ys) - 1)):
        return False, f'三行文本未按"分层任务"→"情境引导"→"跨学科融合"自上而下排列，中心Y={ys}'

    return True, '三行文本分列独立文本框自上而下排列，字号/字体/颜色/居中均满足细则要求'

@rule('D2-39', +3, '右下评价反馈卡片外框')
def _(slide, prs):
    """
    细则：位于距左23.3–29.3cm、距上10.5–15.5cm范围内，宽4–5cm，高3.8–4.6cm，
    整体为右侧较宽、左侧略窄的斜边圆角四边形（parallelogram，水平翻转），
    填充为白色，边线为橙色单实线，线宽1.5–2.2磅。
    检查点（仅按细则约束）：
      1. 位置：left∈[23.3,29.3]cm，top∈[10.5,15.5]cm
      2. 尺寸：宽4–5cm，高3.8–4.6cm
      3. 形状为斜边圆角四边形（parallelogram），且**水平翻转** → 右侧较宽、左侧略窄
      4. 线条转角为圆角 → 圆角表现
      5. 填充为白色
      6. 边线为橙色单实线，线宽1.5–2.2磅
    """
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        if not in_box(l, t, 23.3, 29.3, 10.5, 15.5):
            continue
        if not in_wh(w, h, 4, 5, 3.8, 4.6):
            continue
        if prst_geom(s) != 'parallelogram':
            continue
        if not has_flip_h(s):
            continue
        if not has_round_corners(s):
            continue
        rgb, _ = get_fill_rgb(s)
        if not (rgb is None or is_white_or_near(rgb)):
            continue
        line_rgb = get_line_rgb(s)
        lw = get_line_width_pt(s)
        if not (line_rgb and is_orange(line_rgb)):
            continue
        if not (lw and in_range(lw, 1.5, 2.2)):
            continue
        if has_dash(s):
            continue
        return True, (
            f'{s.name} 右下评价反馈卡片外框 ({l:.2f},{t:.2f}) '
            f'{w:.2f}x{h:.2f}cm parallelogram(水平翻转→右宽左窄) 圆角 '
            f'白色填充 橙色单实线 线宽{lw:.2f}pt'
        )
    return False, (
        '未找到符合以下全部条件的右下评价反馈卡片外框：'
        'left∈[23.3,29.3]cm、top∈[10.5,15.5]cm、宽4–5cm、高3.8–4.6cm、'
        'parallelogram(右宽左窄，水平翻转)、圆角线条转角、白色填充、'
        '橙色单实线边线、线宽1.5–2.2磅'
    )

@rule('D2-40', +1, '右下评价反馈图标')
def _(slide, prs):
    """
    细则：位于右下评价反馈卡片左上部，距左24.6–26.3cm、距上11.0–12.8cm范围内，
    形状为圆形，直径1.3–1.6cm，填充为橙色，内部为白色柱形图或上升趋势图标。
    检查点（仅按细则约束）：
      1. 圆形位置：left∈[24.6,26.3]cm、top∈[11.0,12.8]cm
      2. 直径：宽高均1.3–1.6cm
      3. 填充为橙色
      4. 圆形内部存在白色柱形图或上升趋势图标：圆内至少2个白色可编辑图形/线条/文本部件
    """
    for circle in slide.shapes:
        l, t, w, h = shape_cm(circle)
        if not in_box(l, t, 24.6, 26.3, 11.0, 12.8):
            continue
        if not is_oval_like(circle):
            continue
        if not (in_range(w, 1.3, 1.6) and in_range(h, 1.3, 1.6)):
            continue
        fill_rgb, _ = get_fill_rgb(circle)
        if not (fill_rgb and is_orange(fill_rgb)):
            continue

        white_parts = []
        for part in slide.shapes:
            if part == circle:
                continue
            if not _shape_inside_circle(part, circle, margin=0.1):
                continue
            if _is_white_editable_icon_part(part):
                white_parts.append(part)

        if len(white_parts) >= 2:
            return True, (
                f'{circle.name} 右下评价反馈图标 ({l:.2f},{t:.2f}) '
                f'{w:.2f}x{h:.2f}cm 橙色圆形，内部含{len(white_parts)}个白色柱形图/上升趋势图标部件'
            )
        return False, (
            f'找到橙色圆形图标 {circle.name}，但内部白色可编辑柱形图/上升趋势图标部件仅{len(white_parts)}个（需至少2个）'
        )

    return False, (
        '未找到符合 left∈[24.6,26.3]cm、top∈[11.0,12.8]cm、'
        '直径1.3–1.6cm、橙色填充的圆形图标'
    )

@rule('D2-41', +1, '右下评价反馈标题文本')
def _(slide, prs):
    """
    细则：位于右下评价反馈卡片上部，文本为"评价反馈"，
    字体为黑体、微软雅黑或相近字体，字号16–18磅，加粗，
    颜色为深橙色，位于图标圆右侧。
    检查点（仅按细则约束）：
      1. 文本精确等于"评价反馈"
      2. 位于卡片上部（left∈[26.0,33.0]cm、top∈[11.0,13.8]cm）
      3. 字体为黑体/微软雅黑或相近中文无衬线字体
      4. 字号16–18磅
      5. 加粗
      6. 颜色为深橙色
      7. 位于图标圆右侧（文本left > 图标圆右边缘 - 0.3cm）
    """
    s = find_shape_with_text(slide, '评价反馈')
    if not s:
        return False, '未找到文本精确等于"评价反馈"的形状'

    l, t, w, h = shape_cm(s)
    if not in_box(l, t, 26.0, 33.0, 11.0, 13.8):
        return False, f'"评价反馈"位置({l:.2f},{t:.2f})不在卡片上部left∈[26.0,33.0]、top∈[11.0,13.8]范围内'

    all_runs = _font_info_all_runs(s)
    if not all_runs:
        fi = font_info(s)
        if fi:
            all_runs = [fi]
    if not all_runs:
        return False, '"评价反馈"无法读取字体信息'

    font_ok  = any(is_cjk_sans(r[3]) for r in all_runs)
    sz_ok    = any(r[0] and in_range(r[0], 16, 18) for r in all_runs)
    bold_ok  = any(_is_bold_or_semibold(r[1], r[3]) for r in all_runs)
    color_ok = any(r[2] and (is_deep_orange(r[2]) or is_orange(r[2])) for r in all_runs)

    if not font_ok:
        return False, f'"评价反馈"字体{[r[3] for r in all_runs]}非黑体/微软雅黑或相近字体'
    if not sz_ok:
        return False, f'"评价反馈"字号{[r[0] for r in all_runs]}不在16–18磅范围'
    if not bold_ok:
        return False, '"评价反馈"未加粗'
    if not color_ok:
        return False, f'"评价反馈"颜色{[r[2] for r in all_runs]}非深橙色'

    # 位于图标圆右侧：图标圆在 left∈[24.6,26.3]、top∈[11.0,12.8]，直径1.3–1.6cm
    icon_right = None
    for shape in slide.shapes:
        il, it, iw, ih = shape_cm(shape)
        if (is_oval_like(shape) and
                in_box(il, it, 24.6, 26.3, 11.0, 12.8) and
                in_range(iw, 1.3, 1.6) and in_range(ih, 1.3, 1.6)):
            fill_rgb, _ = get_fill_rgb(shape)
            if fill_rgb and is_orange(fill_rgb):
                icon_right = il + iw
                break

    if icon_right is not None and l < icon_right - 0.3:
        return False, (
            f'"评价反馈"left={l:.2f}cm 不在图标圆右侧（图标右边缘≈{icon_right:.2f}cm）'
        )

    return True, (
        f'{s.name} "评价反馈" 位置({l:.2f},{t:.2f}) '
        f'字号{[r[0] for r in all_runs if r[0]]}pt 加粗 深橙色 图标右侧'
    )

@rule('D2-42', +1, '右下评价反馈分隔线与圆点')
def _(slide, prs):
    """
    细则：位于标题文本下方0.3cm–0.7cm处，为一条浅橙色水平单实线，
    长度3.5cm–4.5cm，线宽0.75–1.25磅；中部有一个橙色小圆点。
    检查点（仅按细则约束）：
      1. 以"评价反馈"标题文本底边为参考，分隔线中心Y在其下方0.3–0.7cm
      2. 分隔线为浅橙色水平单实线，长度3.5–4.5cm，线宽0.75–1.25磅
      3. 分隔线中部有橙色小圆点（圆形，橙色，中心X/Y与线中心偏差≤0.2cm）
    """
    title_shape = find_shape_with_text(slide, '评价反馈')
    if not title_shape:
        return False, '未找到"评价反馈"标题文本作为分隔线参考'

    _, title_top, _, title_h = shape_cm(title_shape)
    title_bottom = title_top + title_h
    lines = []
    dots = []

    for s in slide.shapes:
        line_rgb = get_line_rgb(s)
        line_width = get_line_width_pt(s)
        _, cy = shape_center(s)
        if (is_horizontal_line(s, 3.5, 4.5, 0.12) and
                line_rgb and is_light_orange(line_rgb) and
                line_width and in_range(line_width, 0.75, 1.25) and
                not has_dash(s) and
                in_range(cy - title_bottom, 0.3, 0.7)):
            lines.append(s)

        if is_oval_like(s):
            fill_rgb, _ = get_fill_rgb(s)
            if fill_rgb and is_orange(fill_rgb):
                dots.append(s)

    if not lines:
        return False, '未找到位于"评价反馈"标题下方0.3–0.7cm、长度3.5–4.5cm、浅橙色、线宽0.75–1.25磅的水平单实线'
    if not dots:
        return False, '未找到橙色小圆点'

    for line in lines:
        lx, ly = line_center(line)
        for dot in dots:
            dx, dy = shape_center(dot)
            if abs(dx - lx) <= 0.2 and abs(dy - ly) <= 0.2:
                ll, lt, lw_cm, lh_cm = shape_cm(line)
                return True, (
                    f'{line.name} 位于"评价反馈"下方{ly-title_bottom:.2f}cm，'
                    f'长度{lw_cm:.2f}cm 浅橙色水平单实线；'
                    f'{dot.name} 橙色小圆点位于分隔线中部'
                )

    return False, (
        f'找到分隔线{len(lines)}条、橙色小圆点{len(dots)}个，'
        f'但未满足小圆点位于分隔线中部（中心偏差≤0.2cm）'
    )

@rule('D2-43', +1, '右下评价反馈三行文本')
def _(slide, prs):
    """
    细则：位于右下评价反馈卡片中下部，从上到下出现"过程记录""成果展示""数据改进"，
    字体为黑体、微软雅黑或相近字体，字号13–15磅，颜色为黑色或深灰色，三行均水平居中。
    """
    lo_l, hi_l, lo_t, hi_t = 24, 32, 13.0, 15.5
    required = ('过程记录', '成果展示', '数据改进')

    def line_texts(shape):
        lines = []
        for para in shape.text_frame.paragraphs:
            text = ''.join(r.text for r in para.runs if r.text).strip()
            if text:
                lines.append(text)
        return lines or [shape_text(shape).strip()]

    def all_lines_centered(shape):
        for para in shape.text_frame.paragraphs:
            text = ''.join(r.text for r in para.runs if r.text).strip()
            if not text:
                continue
            align = para.alignment
            ppr = para._p.find('{%s}pPr' % NS['a'])
            xml_center = ppr is not None and ppr.get('algn') == 'ctr'
            if align not in (PP_ALIGN.CENTER, None) and not xml_center:
                return False
        return True

    def style_ok(shape, label):
        all_runs = _font_info_all_runs(shape)
        if not all_runs:
            fi = font_info(shape)
            if fi:
                all_runs = [fi]
        if not all_runs:
            return False, f'"{label}"无法读取字体信息'
        font_ok = any(is_cjk_sans(r[3]) for r in all_runs)
        sz_ok = any(r[0] and in_range(r[0], 13, 15) for r in all_runs)
        color_ok = any(r[2] and _is_dark_or_black(r[2]) for r in all_runs)
        center_ok = all_lines_centered(shape)
        if not font_ok:
            return False, f'"{label}"字体{[r[3] for r in all_runs]}非黑体/微软雅黑或相近字体'
        if not sz_ok:
            return False, f'"{label}"字号{[r[0] for r in all_runs]}不在13–15磅范围'
        if not color_ok:
            return False, f'"{label}"颜色{[r[2] for r in all_runs]}非黑色或深灰色'
        if not center_ok:
            return False, f'"{label}"未水平居中'
        return True, 'ok'

    # 同一文本框三行
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        l, t, w, h = shape_cm(s)
        if not (in_range(l, lo_l, hi_l) and in_range(t, lo_t, hi_t)):
            continue
        if line_texts(s) == list(required):
            ok, reason = style_ok(s, '、'.join(required))
            if not ok:
                return False, reason
            return True, '三行文本在同一文本框内自上而下出现，字号/字体/颜色/居中均满足细则要求'

    # 三个独立文本框
    found = []
    for text in required:
        s = find_text_in_region(slide, text, lo_l, hi_l, lo_t, hi_t)
        if not s:
            return False, f'在右下评价反馈卡片中下部区域未找到文本"{text}"'
        ok, reason = style_ok(s, text)
        if not ok:
            return False, reason
        found.append(s)

    ys = [shape_center(s)[1] for s in found]
    if not all(ys[i] < ys[i + 1] for i in range(len(ys) - 1)):
        return False, f'三行文本未按"过程记录"→"成果展示"→"数据改进"自上而下排列，中心Y={ys}'

    return True, '三行文本分列独立文本框自上而下排列，字号/字体/颜色/居中均满足细则要求'


# ─────────────────── 连接线与箭头 ───────────────────
def connector_rule(slide, region, title):
    lo_l, hi_l, lo_t, hi_t = region
    lines = 0; dots = 0
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        cx, cy = shape_center(s)
        if not (in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t)): continue
        line_rgb = get_line_rgb(s)
        lw = get_line_width_pt(s)
        if prst_geom(s) == 'line' and line_rgb and is_orange(line_rgb) and lw and in_range(lw, 1.2, 1.8):
            lines += 1
        if is_oval_like(s) and in_range(w, 0.18, 0.32) and in_range(h, 0.18, 0.32):
            fill_rgb, _ = get_fill_rgb(s)
            if fill_rgb and is_orange(fill_rgb):
                dots += 1
    if lines >= 1 and dots >= 2:
        return True, f'{title}检测到{lines}条橙色折线/线段和{dots}个橙色圆点'
    return False, f'{title}线段{lines} 圆点{dots}'

def _connector_points(shape):
    l, t, w, h = shape_cm(shape)
    return [
        (l, t), (l + w, t), (l, t + h), (l + w, t + h),
        (l + w / 2, t), (l + w / 2, t + h), (l, t + h / 2), (l + w, t + h / 2),
        (l + w / 2, t + h / 2),
    ]


def _dot_near_connector(dot, connector, max_dist=0.35):
    dx, dy = shape_center(dot)
    for px, py in _connector_points(connector):
        if abs(dx - px) <= max_dist and abs(dy - py) <= max_dist:
            return True
    return False


def _orange_solid_dot(shape, lo=0.18, hi=0.32):
    if not is_oval_like(shape):
        return False
    _, _, w, h = shape_cm(shape)
    if not (in_range(w, lo, hi) and in_range(h, lo, hi)):
        return False
    fill_rgb, _ = get_fill_rgb(shape)
    return bool(fill_rgb and is_orange(fill_rgb))


def _orange_polyline_connector(shape, lo_l, hi_l, lo_t, hi_t):
    l, t, w, h = shape_cm(shape)
    cx, cy = shape_center(shape)
    if not (in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t)):
        return False
    if not _is_office_line_or_connector(shape):
        return False
    line_rgb = get_line_rgb(shape)
    lw = get_line_width_pt(shape)
    if not (line_rgb and is_orange(line_rgb)):
        return False
    if not (lw and in_range(lw, 1.2, 1.8)):
        return False
    return True


@rule('D2-44', +1, '左上连接折线与圆点')
def _(slide, prs):
    """
    细则：位于左上资源整合卡片右侧与中部椭圆左上之间，距左9–12.5cm、距上5.2–7.5cm范围内，
    为一条橙色折线，线宽1.2–1.8磅；折线两端或转折处有2个橙色实心圆点，直径0.18–0.32cm。
    检查点（仅按细则约束）：
      1. 折线/连接符中心位于 left∈[9,12.5]cm、top∈[5.2,7.5]cm区域
      2. 折线为橙色，线宽1.2–1.8磅
      3. 区域内有2个橙色实心圆点，直径0.18–0.32cm
      4. 两个圆点位于折线端点或转折处附近（与折线边界关键点偏差≤0.35cm）
    """
    lo_l, hi_l, lo_t, hi_t = 9, 12.5, 5.2, 7.5
    connectors = [s for s in slide.shapes if _orange_polyline_connector(s, lo_l, hi_l, lo_t, hi_t)]
    dots = []
    for s in slide.shapes:
        cx, cy = shape_center(s)
        if not (in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t)):
            continue
        if _orange_solid_dot(s, 0.18, 0.32):
            dots.append(s)

    if not connectors:
        return False, '未找到位于left∈[9,12.5]cm、top∈[5.2,7.5]cm区域内、橙色、线宽1.2–1.8磅的折线/连接符'
    if len(dots) < 2:
        return False, f'橙色实心圆点仅{len(dots)}个，需2个且直径0.18–0.32cm'

    for conn in connectors:
        near_dots = [d for d in dots if _dot_near_connector(d, conn, max_dist=0.35)]
        if len(near_dots) >= 2:
            l, t, w, h = shape_cm(conn)
            lw = get_line_width_pt(conn)
            return True, (
                f'{conn.name} 左上连接折线 ({l:.2f},{t:.2f}) {w:.2f}x{h:.2f}cm '
                f'橙色线宽{lw:.2f}pt，折线两端/转折处有{len(near_dots)}个橙色实心圆点'
            )

    return False, f'找到折线{len(connectors)}条、橙色实心圆点{len(dots)}个，但不足2个圆点位于折线两端或转折处'

@rule('D2-45', +1, '左下连接折线与圆点')
def _(slide, prs):
    """
    细则：位于左下活动设计卡片右侧与中部椭圆左下之间，距左9.5–13cm、距上11.0–13.5cm范围内，
    为一条橙色折线，线宽1.2–1.8磅；折线两端或转折处有2个橙色实心圆点，直径0.18–0.32cm。
    检查点（仅按细则约束）：
      1. 折线/连接符中心位于 left∈[9.5,13]cm、top∈[11.0,13.5]cm区域
      2. 折线为橙色，线宽1.2–1.8磅
      3. 区域内有2个橙色实心圆点，直径0.18–0.32cm
      4. 两个圆点位于折线端点或转折处附近（与折线边界关键点偏差≤0.35cm）
    """
    lo_l, hi_l, lo_t, hi_t = 9.5, 13, 11.0, 13.5
    connectors = [s for s in slide.shapes if _orange_polyline_connector(s, lo_l, hi_l, lo_t, hi_t)]
    dots = []
    for s in slide.shapes:
        cx, cy = shape_center(s)
        if not (in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t)):
            continue
        if _orange_solid_dot(s, 0.18, 0.32):
            dots.append(s)

    if not connectors:
        return False, '未找到位于left∈[9.5,13]cm、top∈[11.0,13.5]cm区域内、橙色、线宽1.2–1.8磅的折线/连接符'
    if len(dots) < 2:
        return False, f'橙色实心圆点仅{len(dots)}个，需2个且直径0.18–0.32cm'

    for conn in connectors:
        near_dots = [d for d in dots if _dot_near_connector(d, conn, max_dist=0.35)]
        if len(near_dots) >= 2:
            l, t, w, h = shape_cm(conn)
            lw = get_line_width_pt(conn)
            return True, (
                f'{conn.name} 左下连接折线 ({l:.2f},{t:.2f}) {w:.2f}x{h:.2f}cm '
                f'橙色线宽{lw:.2f}pt，折线两端/转折处有{len(near_dots)}个橙色实心圆点'
            )

    return False, f'找到折线{len(connectors)}条、橙色实心圆点{len(dots)}个，但不足2个圆点位于折线两端或转折处'

@rule('D2-46', +1, '右上连接折线与圆点')
def _(slide, prs):
    """
    细则：位于右上教学支持卡片左侧与中部椭圆右上之间，距左21.0–24.1cm、距上5.2–7.5cm范围内，
    为一条橙色折线，线宽1.2–1.8磅；折线两端或转折处有2个橙色实心圆点，直径0.18–0.32cm。
    检查点（仅按细则约束）：
      1. 折线/连接符中心位于 left∈[21.0,24.1]cm、top∈[5.2,7.5]cm区域
      2. 折线为橙色，线宽1.2–1.8磅
      3. 区域内有2个橙色实心圆点，直径0.18–0.32cm
      4. 两个圆点位于折线端点或转折处附近（与折线边界关键点偏差≤0.35cm）
    """
    lo_l, hi_l, lo_t, hi_t = 21.0, 24.1, 5.2, 7.5
    connectors = [s for s in slide.shapes if _orange_polyline_connector(s, lo_l, hi_l, lo_t, hi_t)]
    dots = []
    for s in slide.shapes:
        cx, cy = shape_center(s)
        if not (in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t)):
            continue
        if _orange_solid_dot(s, 0.18, 0.32):
            dots.append(s)

    if not connectors:
        return False, '未找到位于left∈[21.0,24.1]cm、top∈[5.2,7.5]cm区域内、橙色、线宽1.2–1.8磅的折线/连接符'
    if len(dots) < 2:
        return False, f'橙色实心圆点仅{len(dots)}个，需2个且直径0.18–0.32cm'

    for conn in connectors:
        near_dots = [d for d in dots if _dot_near_connector(d, conn, max_dist=0.35)]
        if len(near_dots) >= 2:
            l, t, w, h = shape_cm(conn)
            lw = get_line_width_pt(conn)
            return True, (
                f'{conn.name} 右上连接折线 ({l:.2f},{t:.2f}) {w:.2f}x{h:.2f}cm '
                f'橙色线宽{lw:.2f}pt，折线两端/转折处有{len(near_dots)}个橙色实心圆点'
            )

    return False, f'找到折线{len(connectors)}条、橙色实心圆点{len(dots)}个，但不足2个圆点位于折线两端或转折处'

@rule('D2-47', +1, '右下连接折线与圆点')
def _(slide, prs):
    """
    细则：位于右下评价反馈卡片左侧与中部椭圆右下之间，距左21.5–24.5cm、距上11.0–13.5cm范围内，
    为一条橙色折线，线宽1.2–1.8磅；折线两端或转折处有2个橙色实心圆点，直径0.18–0.32cm。
    检查点（仅按细则约束）：
      1. 折线/连接符中心位于 left∈[21.5,24.5]cm、top∈[11.0,13.5]cm区域
      2. 折线为橙色，线宽1.2–1.8磅
      3. 区域内有2个橙色实心圆点，直径0.18–0.32cm
      4. 两个圆点位于折线端点或转折处附近（与折线边界关键点偏差≤0.35cm）
    """
    lo_l, hi_l, lo_t, hi_t = 21.5, 24.5, 11.0, 13.5
    connectors = [s for s in slide.shapes if _orange_polyline_connector(s, lo_l, hi_l, lo_t, hi_t)]
    dots = []
    for s in slide.shapes:
        cx, cy = shape_center(s)
        if not (in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t)):
            continue
        if _orange_solid_dot(s, 0.18, 0.32):
            dots.append(s)

    if not connectors:
        return False, '未找到位于left∈[21.5,24.5]cm、top∈[11.0,13.5]cm区域内、橙色、线宽1.2–1.8磅的折线/连接符'
    if len(dots) < 2:
        return False, f'橙色实心圆点仅{len(dots)}个，需2个且直径0.18–0.32cm'

    for conn in connectors:
        near_dots = [d for d in dots if _dot_near_connector(d, conn, max_dist=0.35)]
        if len(near_dots) >= 2:
            l, t, w, h = shape_cm(conn)
            lw = get_line_width_pt(conn)
            return True, (
                f'{conn.name} 右下连接折线 ({l:.2f},{t:.2f}) {w:.2f}x{h:.2f}cm '
                f'橙色线宽{lw:.2f}pt，折线两端/转折处有{len(near_dots)}个橙色实心圆点'
            )

    return False, f'找到折线{len(connectors)}条、橙色实心圆点{len(dots)}个，但不足2个圆点位于折线两端或转折处'

def _has_arrows_both_ends(shape):
    el = shape._element
    head = el.find('.//{%s}headEnd' % NS['a'])
    tail = el.find('.//{%s}tailEnd' % NS['a'])
    def end_has_arrow(end):
        if end is None:
            return False
        return end.get('type', 'none') not in ('none', '')
    return end_has_arrow(head) and end_has_arrow(tail)


def _is_integrated_double_arrow(shape, orientation):
    """一体式双向箭头形状：Office/WPS 预设的整块双向箭头
    （leftRightArrow / upDownArrow 及其 callout 变体）。
    这类形状本身的几何即包含两端箭头，无需再检查 headEnd/tailEnd。
    orientation: 'horizontal' 或 'vertical'。
    """
    g = prst_geom(shape)
    if orientation == 'horizontal':
        return g in ('leftRightArrow', 'leftRightArrowCallout')
    if orientation == 'vertical':
        return g in ('upDownArrow', 'upDownArrowCallout')
    return False


def _integrated_arrow_color(shape):
    """一体式箭头颜色：优先取实/渐变填充主色，回退到线条色。"""
    rgb, ft = get_fill_rgb(shape)
    if ft in ('solid', 'gradient') and rgb:
        return rgb
    return get_line_rgb(shape)


@rule('D2-48', +1, '左右两侧上下卡片之间竖向虚线双箭头')
def _(slide, prs):
    """
    细则：位于距左6.0–6.6cm、距上7.2–10.8cm和距左26.6–27.2cm、距上7.2–10.8cm范围内，
    为橙色竖向虚线双箭头，高度2.7–3.3cm，线宽1–1.5磅，上下两端均有箭头。
    检查点（仅按细则约束）：
      1. 左右两个区域各有一条竖向线/连接符 或 一体式竖向双向箭头
      2. 位置中心在对应区域内
      3. 高度2.7–3.3cm，且高度大于宽度（竖向）
      4. 橙色，虚线视觉，线宽1–1.5磅（**一体式双向箭头同样要求**：
         其填充需为 noFill/无实心（避免实心块），线条须为虚线/点划线/系统虚线，且线宽1–1.5磅）
      5. 上下两端均有箭头（一体式双向箭头本身即满足）
    """
    hits = []
    for region in [(6.0, 6.6, 7.2, 10.8), (26.6, 27.2, 7.2, 10.8)]:
        lo_l, hi_l, lo_t, hi_t = region
        found = None
        for s in slide.shapes:
            cx, cy = shape_center(s)
            if not (in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t)):
                continue
            _, _, w, h = shape_cm(s)
            if not (in_range(h, 2.7, 3.3) and h > w):
                continue
            # 分支 A：一体式竖向双向箭头（upDownArrow / callout）——形状本身即代表双箭头
            #        但仍需按 rubric 强制"虚线视觉 + 线宽1–1.5磅"：
            #        (a) 填充必须为 noFill / 无填充节点（**白名单**，避免 solid/gradient/
            #            pattern/blip 等任何视觉上呈实心块的填充通过）；
            #        (b) 线条为虚线/点划线（has_dash）；
            #        (c) 线宽在 1–1.5磅；
            #        (d) 线条为橙色。
            if _is_integrated_double_arrow(s, 'vertical'):
                # (a) 只允许 无填充 / noFill：显式 <a:noFill> 或 spPr 中根本无填充节点
                sp_pr = s._element.find('.//{%s}spPr' % NS['a'])
                has_no_fill_node = (
                    sp_pr is not None and
                    sp_pr.find('{%s}noFill' % NS['a']) is not None
                )
                any_fill_node = False
                if sp_pr is not None:
                    for _tag in ('solidFill', 'gradFill', 'pattFill',
                                 'blipFill', 'grpFill'):
                        if sp_pr.find('{%s}%s' % (NS['a'], _tag)) is not None:
                            any_fill_node = True
                            break
                if not (has_no_fill_node or (sp_pr is not None and not any_fill_node)):
                    continue
                # (b)(c) 线条虚线 & 线宽 1–1.5磅
                if not has_dash(s):
                    continue
                lw = get_line_width_pt(s)
                if not (lw and in_range(lw, 1, 1.5)):
                    continue
                # (d) 线条橙色
                line_rgb = get_line_rgb(s)
                if not (line_rgb and is_orange(line_rgb)):
                    continue
                found = s
                break
            # 分支 B：line / 连接符 拼合的虚线双箭头
            if not _is_office_line_or_connector(s):
                continue
            line_rgb = get_line_rgb(s)
            lw = get_line_width_pt(s)
            if not (line_rgb and is_orange(line_rgb)):
                continue
            if not (lw and in_range(lw, 1, 1.5)):
                continue
            if not has_dash(s):
                continue
            if not _has_arrows_both_ends(s):
                continue
            found = s
            break
        if found:
            hits.append(found)

    if len(hits) == 2:
        return True, (
            '左右两侧均检测到橙色竖向虚线双箭头（一体式或连接符拼合），'
            '高度2.7–3.3cm、虚线线宽1–1.5磅'
        )
    return False, (
        f'检测到{len(hits)}/2个符合区域、高度、橙色、虚线且线宽1–1.5磅的竖向双箭头'
    )

@rule('D2-49', +1, '椭圆左右两侧中部水平灰色双向箭头')
def _(slide, prs):
    """
    细则："实施枢纽"外的椭圆左右两侧中部水平灰色双向箭头，
    位于距左6.5–9.5cm、距上8.5–9.5cm和距左24–27cm、距上8.5–9.5cm，
    为灰色水平虚线或点划线双向箭头，长度1.8–2.2cm，线宽1–1.5磅，左右两端均有箭头。
    检查点（仅按细则约束）：
      1. 左右两个区域各有一条水平线/连接符 或 一体式水平双向箭头
      2. 位置中心在对应区域内
      3. 长度1.8–2.2cm，且宽度大于高度（水平）
      4. 灰色，虚线或点划线，线宽1–1.5磅（**一体式双向箭头同样要求**：
         其填充需为 noFill/无实心（避免实心块），线条须为虚线/点划线，且线宽1–1.5磅）
      5. 左右两端均有箭头（一体式双向箭头本身即满足）
    """
    hits = []
    for region in [(6.5, 9.5, 8.5, 9.5), (24.0, 27.0, 8.5, 9.5)]:
        lo_l, hi_l, lo_t, hi_t = region
        found = None
        for s in slide.shapes:
            cx, cy = shape_center(s)
            if not (in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t)):
                continue
            _, _, w, h = shape_cm(s)
            if not (in_range(w, 1.8, 2.2) and w > h):
                continue
            # 分支 A：一体式水平双向箭头
            #        按 rubric 强制"虚线/点划线视觉 + 线宽1–1.5磅"：
            #        (a) 填充必须为 noFill / 无填充节点（**白名单**，避免 solid/gradient/
            #            pattern/blip 等任何视觉上呈实心块的填充通过——一体式实心箭头
            #            无法在视觉上表达"虚线/点划线"，rubric 建议直接不接受）；
            #        (b) 线条为虚线/点划线（has_dash）；
            #        (c) 线宽在 1–1.5磅；
            #        (d) 线条为灰色。
            if _is_integrated_double_arrow(s, 'horizontal'):
                # (a) 只允许 无填充 / noFill
                sp_pr = s._element.find('.//{%s}spPr' % NS['a'])
                has_no_fill_node = (
                    sp_pr is not None and
                    sp_pr.find('{%s}noFill' % NS['a']) is not None
                )
                any_fill_node = False
                if sp_pr is not None:
                    for _tag in ('solidFill', 'gradFill', 'pattFill',
                                 'blipFill', 'grpFill'):
                        if sp_pr.find('{%s}%s' % (NS['a'], _tag)) is not None:
                            any_fill_node = True
                            break
                if not (has_no_fill_node or (sp_pr is not None and not any_fill_node)):
                    continue
                if not has_dash(s):
                    continue
                lw = get_line_width_pt(s)
                if not (lw and in_range(lw, 1, 1.5)):
                    continue
                line_rgb = get_line_rgb(s)
                if not (line_rgb and is_gray(line_rgb)):
                    continue
                found = s
                break
            # 分支 B：line / 连接符拼合的水平双向箭头
            if not _is_office_line_or_connector(s):
                continue
            line_rgb = get_line_rgb(s)
            lw = get_line_width_pt(s)
            if not (line_rgb and is_gray(line_rgb)):
                continue
            if not (lw and in_range(lw, 1, 1.5)):
                continue
            if not has_dash(s):
                continue
            if not _has_arrows_both_ends(s):
                continue
            found = s
            break
        if found:
            hits.append(found)

    if len(hits) == 2:
        return True, (
            '椭圆左右两侧中部均检测到灰色水平虚线双向箭头（一体式或连接符拼合），'
            '长度1.8–2.2cm、虚线线宽1–1.5磅'
        )
    return False, (
        f'检测到{len(hits)}/2个符合区域、长度、灰色、虚线/点划线且线宽1–1.5磅的水平双向箭头'
    )

def _compute_d2_50_regions(slide):
    """基于"成效层"箭头(下方 downArrow)和下方两卡片边框，动态计算左右连接区域。

    - 左连接区域: X 位于"左下卡片右边"与"成效层箭头左边"之间(取二者相对的向内区间);
                  Y 位于成效层箭头竖向中点 ± 容差.
    - 右连接区域: 对称计算.

    若任一锚点缺失, 返回宽松兜底区域(与原硬编码区间近似, 保证不比原代码更严).
    返回: [(lo_l, hi_l, lo_t, hi_t), (lo_l, hi_l, lo_t, hi_t)]  # [left, right]
    """
    fallback = [(3.5, 7.5, 14.5, 15.5), (21.5, 25.8, 14.5, 15.5)]

    arrow = find_bottom_arrow_body(slide)
    if arrow is None:
        return fallback
    al, at, aw, ah = shape_cm(arrow)
    a_left, a_right = al, al + aw
    a_cy = at + ah / 2.0

    # 下方左卡片(D2-29 特征): parallelogram 未翻转, 在其位置/尺寸区间内
    left_card = None
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        if not in_box(l, t, 4.3, 10.8, 10.5, 15.5):
            continue
        if not in_wh(w, h, 4, 5, 3.8, 4.6):
            continue
        if prst_geom(s) != 'parallelogram':
            continue
        if has_flip_h(s):
            continue
        left_card = s
        break

    # 下方右卡片(D2-39 特征): parallelogram 水平翻转
    right_card = None
    for s in slide.shapes:
        l, t, w, h = shape_cm(s)
        if not in_box(l, t, 23.3, 29.3, 10.5, 15.5):
            continue
        if not in_wh(w, h, 4, 5, 3.8, 4.6):
            continue
        if prst_geom(s) != 'parallelogram':
            continue
        if not has_flip_h(s):
            continue
        right_card = s
        break

    # Y 带: 成效层竖向中点 ± 1.0cm
    y_tol = 1.0
    lo_t = a_cy - y_tol
    hi_t = a_cy + y_tol

    # X 带: 左连接 = [左卡片右边 - 0.5, 成效层左边 + 0.5]; 右连接 = [成效层右边 - 0.5, 右卡片左边 + 0.5]
    if left_card is not None:
        lcl, _lct, lcw, _lch = shape_cm(left_card)
        left_card_right = lcl + lcw
        left_lo_l = min(left_card_right, a_left) - 0.5
        left_hi_l = max(left_card_right, a_left) + 0.5
    else:
        left_lo_l = max(0.0, a_left - 4.5)
        left_hi_l = a_left + 0.5

    if right_card is not None:
        rcl, _rct, _rcw, _rch = shape_cm(right_card)
        right_card_left = rcl
        right_lo_l = min(a_right, right_card_left) - 0.5
        right_hi_l = max(a_right, right_card_left) + 0.5
    else:
        right_lo_l = a_right - 0.5
        right_hi_l = a_right + 4.5

    left_region = (left_lo_l, left_hi_l, lo_t, hi_t)
    right_region = (right_lo_l, right_hi_l, lo_t, hi_t)

    # 与原硬编码兜底区间取并集, 保证动态计算不比原代码更严, 避免把老素材排除.
    def _union(dyn, fb):
        return (min(dyn[0], fb[0]), max(dyn[1], fb[1]),
                min(dyn[2], fb[2]), max(dyn[3], fb[3]))

    return [_union(left_region, fallback[0]), _union(right_region, fallback[1])]


@rule('D2-50', +1, '成效层左右到下方卡片间水平灰色双向箭头')
def _(slide, prs):
    """
    细则："成效层"箭头卡片边框左右两侧到下方两卡片间各有一个水平的灰色双向箭头，
    长度3.3–3.7cm，为灰色水平虚线或点划线双向箭头，线宽1–1.5磅，左右两端均有箭头。
    检查点（仅按细则约束）：
      1. 左右两个下方连接区域各有一条水平线/连接符 或 一体式水平双向箭头。
         区域由"成效层"箭头(find_bottom_arrow_body) 与下方两卡片(D2-29 未翻转 /
         D2-39 水平翻转的 parallelogram) 几何动态计算，而非固定硬编码坐标。
      2. 长度3.3–3.7cm，且宽度大于高度（水平）
      3. 灰色，虚线或点划线，线宽1–1.5磅
         **一体式双向箭头同样要求**（对齐 D2-48/D2-49 处理）:
           (a) 填充必须为 noFill / 无填充节点（白名单：`spPr` 显式 `<a:noFill>` 或
               `spPr` 中无任何填充节点；`solidFill`/`gradFill`/`pattFill`/`blipFill`/
               `grpFill` 一律拒——实心一体式箭头视觉上无法表达虚线/点划线，不接受)；
           (b) 线条为虚线/点划线(has_dash)；
           (c) 线宽 1–1.5磅；
           (d) 线条颜色为灰色.
      4. 左右两端均有箭头（一体式双向箭头本身即满足）
    """
    regions = _compute_d2_50_regions(slide)
    hits = []
    for region in regions:
        lo_l, hi_l, lo_t, hi_t = region
        found = None
        for s in slide.shapes:
            cx, cy = shape_center(s)
            if not (in_range(cx, lo_l, hi_l) and in_range(cy, lo_t, hi_t)):
                continue
            _, _, w, h = shape_cm(s)
            if not (in_range(w, 3.3, 3.7) and w > h):
                continue
            # 分支 A：一体式水平双向箭头（与 D2-48/D2-49 同口径）
            if _is_integrated_double_arrow(s, 'horizontal'):
                # (a) 只允许 无填充 / noFill
                sp_pr = s._element.find('.//{%s}spPr' % NS['a'])
                has_no_fill_node = (
                    sp_pr is not None and
                    sp_pr.find('{%s}noFill' % NS['a']) is not None
                )
                any_fill_node = False
                if sp_pr is not None:
                    for _tag in ('solidFill', 'gradFill', 'pattFill',
                                 'blipFill', 'grpFill'):
                        if sp_pr.find('{%s}%s' % (NS['a'], _tag)) is not None:
                            any_fill_node = True
                            break
                if not (has_no_fill_node or (sp_pr is not None and not any_fill_node)):
                    continue
                if not has_dash(s):
                    continue
                lw = get_line_width_pt(s)
                if not (lw and in_range(lw, 1, 1.5)):
                    continue
                line_rgb = get_line_rgb(s)
                if not (line_rgb and is_gray(line_rgb)):
                    continue
                found = s
                break
            # 分支 B：line / 连接符拼合的水平双向箭头
            if not _is_office_line_or_connector(s):
                continue
            line_rgb = get_line_rgb(s)
            lw = get_line_width_pt(s)
            if not (line_rgb and is_gray(line_rgb)):
                continue
            if not (lw and in_range(lw, 1, 1.5)):
                continue
            if not has_dash(s):
                continue
            if not _has_arrows_both_ends(s):
                continue
            found = s
            break
        if found:
            hits.append(found)

    if len(hits) == 2:
        return True, (
            '成效层左右到下方两卡片间均检测到灰色水平双向箭头(一体式或虚线双箭头连接符)，'
            '长度3.3–3.7cm、虚线/点划线线宽1–1.5磅；'
            '连接区域由成效层箭头+下方两卡片几何动态计算'
        )
    return False, (
        f'检测到{len(hits)}/2个符合区域、长度、灰色、虚线/点划线且线宽1–1.5磅的成效层水平双向箭头'
    )

def _normalize_caption_text(text):
    return re.sub(r'\s+', ' ', text).strip()


def _is_black(rgb):
    if rgb is None:
        return False
    return max(rgb) <= 40


def _is_song_or_hei_or_similar(name):
    return is_cjk_sans(name) or is_cjk_serif(name)


@rule('D2-51', +1, '底部图注编号文本')
def _(slide, prs):
    """
    细则：位于距左10.5–23.5cm、距上18–19cm范围内，
    文本为"图3-2 数字化校园阅读项目实施的三维协同模型"，
    字体为宋体、黑体或相近中文字体，字号14–16磅，颜色为黑色，加粗或半加粗。
    检查点（仅按细则约束）：
      1. 文本规范化后精确等于指定图注
      2. 位置：left∈[10.5,23.5]cm，top∈[18,19]cm
      3. 字体为宋体、黑体或相近中文字体
      4. 字号14–16磅
      5. 颜色为黑色
      6. 加粗或半加粗
    """
    expected = '图3-2 数字化校园阅读项目实施的三维协同模型'
    for s in slide.shapes:
        if not s.has_text_frame:
            continue
        txt = _normalize_caption_text(shape_text(s))
        if txt != expected:
            continue
        l, t, w, h = shape_cm(s)
        if not in_box(l, t, 10.5, 23.5, 18.0, 19.0):
            return False, f'图注位置({l:.2f},{t:.2f})不在left∈[10.5,23.5]cm、top∈[18,19]cm范围内'

        all_runs = _font_info_all_runs(s)
        if not all_runs:
            fi = font_info(s)
            if fi:
                all_runs = [fi]
        if not all_runs:
            return False, '图注无法读取字体信息'

        font_ok = any(_is_song_or_hei_or_similar(r[3]) for r in all_runs)
        sz_ok = any(r[0] and in_range(r[0], 14, 16) for r in all_runs)
        color_ok = any(r[2] and _is_black(r[2]) for r in all_runs)
        bold_ok = any(_is_bold_or_semibold(r[1], r[3]) for r in all_runs)

        if not font_ok:
            return False, f'图注字体{[r[3] for r in all_runs]}非宋体、黑体或相近中文字体'
        if not sz_ok:
            return False, f'图注字号{[r[0] for r in all_runs]}不在14–16磅范围'
        if not color_ok:
            return False, f'图注颜色{[r[2] for r in all_runs]}非黑色'
        if not bold_ok:
            return False, '图注未加粗或半加粗'

        return True, (
            f'{s.name} 图注位于({l:.2f},{t:.2f}) '
            f'字号{[r[0] for r in all_runs if r[0]]}pt 黑色 加粗/半加粗'
        )
    return False, '未找到文本精确为"图3-2 数字化校园阅读项目实施的三维协同模型"的底部图注文本'

# ─────────────────── D2-52 整体排版规则已按用户指示删除 ───────────────────
# 原 D2-52 及其独占 helpers(_rule_passed / _shape_within_slide / _shape_overlap_area /
# _shape_area / _text_exceeds_frame) 一并移除, 不再对"主要对象命中+可视范围+文本超框+
# 图形重叠"做整体排版判定。_find_impl_frame 由其它规则(D2-11 等) 复用, 保留。





# ─────────────────── 运行与输出 ───────────────────
def _locate_document(dir_path):
    """在给定目录里定位被评估的 .pptx 文件。
    优先匹配预期文件名；否则按扩展名扫描目录，返回第一个匹配项。
    仅识别 .pptx (不再回退到 .ppt); 找不到时返回 None。"""
    if not dir_path or not os.path.isdir(dir_path):
        return None
    preferred = os.path.join(dir_path, PREFERRED_DOC_NAME)
    if os.path.isfile(preferred):
        return preferred
    for name in sorted(os.listdir(dir_path)):
        if name.startswith('~$'):
            continue
        if name.lower().endswith('.pptx'):
            return os.path.join(dir_path, name)
    return None


# ─────────────────── detail 语义识别 ───────────────────
# 说明：按项目要求，对外返回结构中的 detail 字段一律置空，不再输出规则的
# 通过/跳过/异常原因；此处不再保留相关识别函数，评分逻辑与输出结构不变。


def _build_result(file_name, status='ok', error=None,
                  dim1_pass=False, dim1_reason='',
                  dim2_items=None, total_score=0, max_score=0):
    return {
        'id': SCRIPT_ID,
        'file_name': file_name or '',
        'status': status,
        'error': error,
        'dim1_pass': dim1_pass,
        'dim1_reason': dim1_reason,
        'dim2_items': dim2_items or [],
        'total_score': total_score,
        'max_score': max_score,
    }


def evaluate(dir_path: str) -> dict:
    """按统一约定评估：接收脚本所在目录路径，脚本自己在目录里定位被评估文档。

    返回结构见"脚本接口差异与统一建议.md" §2.2。
    """
    max_score = sum(score for _, score, _, _ in RULES)

    try:
        filepath = _locate_document(dir_path)
        if not filepath:
            return _build_result(
                file_name='',
                status='error',
                error=f'目录 {dir_path!r} 中未找到 .pptx 文件',
                max_score=max_score,
            )
        file_name = os.path.basename(filepath)

        try:
            prs = Presentation(filepath)
            invalid_coordinates = normalize_integral_coordinates(prs)
        except Exception as e:
            return _build_result(
                file_name=file_name,
                status='error',
                error=f'文件无法正常打开：{type(e).__name__}: {e}',
                max_score=max_score,
            )

        if invalid_coordinates:
            samples = ', '.join(invalid_coordinates[:3])
            return _build_result(
                file_name=file_name,
                status='ok',
                dim1_pass=False,
                dim1_reason=(
                    f'PPT XML 包含非整数坐标（共 {len(invalid_coordinates)} 处，'
                    f'示例：{samples}）'
                ),
                max_score=max_score,
            )

        dim1_ok, reasons = check_dim1(prs, filepath)
        if not dim1_ok:
            return _build_result(
                file_name=file_name,
                status='ok',
                dim1_pass=False,
                dim1_reason='；'.join(reasons),
                dim2_items=[],
                total_score=0,
                max_score=max_score,
            )

        slide = prs.slides[0]
        dim2_items = []
        total_score = 0
        for rid, score, title, checker in RULES:
            ok = False
            try:
                result = checker(slide, prs)
                if isinstance(result, tuple) and len(result) >= 1:
                    ok = bool(result[0])
            except Exception:
                # 单条规则执行异常按未通过处理，避免影响其它规则的评分。
                ok = False

            # 按要求：对外输出的 detail 字段一律置空，不改变评分逻辑与其它输出结构。
            dim2_items.append({
                'rule': f'{rid} {title}',
                'max_delta': score,
                'delta': score if ok else 0,
                'hit': bool(ok),
                'detail': '',
            })
            if ok:
                total_score += score

        return _build_result(
            file_name=file_name,
            status='ok',
            dim1_pass=True,
            dim1_reason='',
            dim2_items=dim2_items,
            total_score=total_score,
            max_score=max_score,
        )
    except Exception as e:
        return _build_result(
            file_name=os.path.basename(dir_path or ''),
            status='error',
            error=f'{type(e).__name__}: {e}',
            max_score=max_score,
        )


if __name__ == '__main__':
    # 本地调试入口：默认使用脚本所在目录，可通过命令行覆盖
    _target_dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(_target_dir), ensure_ascii=False, indent=2))
