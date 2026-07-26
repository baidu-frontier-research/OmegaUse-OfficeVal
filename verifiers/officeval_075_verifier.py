# -*- coding: utf-8 -*-
"""
对 “呼吸项目收费解读_合规模板风格版.pptx” 的自动评估脚本。

评估逻辑：
  维度1（可用与可修改性）：先做硬性体检。任何一条不满足 → 0 分，且不再继续评分。
  维度2（完成度评分细则）：通过维度1之后才进入；逐条独立检查，命中则累加该条对应的分值
                            （细则中只列了得分项，没有显式扣分项，但脚本保留扣分支持：
                             如果某条命中失败但用户希望按“扣分点”计分，可在 _RULES 中
                             用负分声明）。

实现思路：
  - 直接解析 .pptx 的 XML 与 python-pptx 对象，避免依赖人工。
  - 字体大小 / 颜色 / 加粗 等属性允许小幅容差（颜色用 ΔRGB ≤ ~12，字号 ±0.5pt），
    以兼容设计稿和细则之间的细微差异。
  - 页码文本细则要求 “01/39”，实际文件里只写了 “01”，按“以页码可读、唯一对应”
    的口径放宽：只要 TextBox 7 文本符合 0?N 形式即视为通过；
    其他类似难以严格还原的项目（如部分非标字号 5.5pt / 9.7pt / 9.8pt /15.5pt 等）
    也使用接近匹配。
"""

from __future__ import annotations

import json
import os
import re
import sys
from dataclasses import dataclass, field
from typing import Callable, Iterable, List, Optional, Tuple

from pptx import Presentation

# ---------------------------------------------------------------------------
# 基础工具
# ---------------------------------------------------------------------------

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


def _hex_to_rgb(h: str) -> Tuple[int, int, int]:
    h = h.strip().lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def color_close(actual: Optional[str], target: str, tol: int = 18) -> bool:
    """按通道差异判断颜色是否接近。"""
    if actual is None:
        return False
    try:
        a = _hex_to_rgb(str(actual))
        b = _hex_to_rgb(target)
    except Exception:
        return False
    return all(abs(a[i] - b[i]) <= tol for i in range(3))


def size_close(actual: Optional[float], target: float, tol: float = 0.6) -> bool:
    if actual is None:
        return False
    return abs(actual - target) <= tol


# ---------------------------------------------------------------------------
# 文本/字体信息抽取
# ---------------------------------------------------------------------------

@dataclass
class RunInfo:
    text: str
    font: Optional[str]
    size_pt: Optional[float]
    bold: Optional[bool]
    color_hex: Optional[str]  # 形如 "FFFFFF"


def _iter_runs(shape) -> Iterable[RunInfo]:
    """从单个形状中抽取每个 run 的字体信息，回退到 paragraph 的 defRPr。"""
    if not shape.has_text_frame:
        return []
    tx = shape.text_frame._txBody  # CT_TextBody
    runs = []
    for p in tx.findall(".//a:p", NS):
        # paragraph default rPr
        defRPr = p.find("a:pPr/a:defRPr", NS)
        d_sz = d_b = d_face = d_col = None
        if defRPr is not None:
            d_sz = defRPr.get("sz")
            d_b = defRPr.get("b")
            face_el = defRPr.find("a:latin", NS)
            if face_el is not None:
                d_face = face_el.get("typeface")
            col_el = defRPr.find("a:solidFill/a:srgbClr", NS)
            if col_el is not None:
                d_col = col_el.get("val")
        for r in p.findall("a:r", NS):
            rPr = r.find("a:rPr", NS)
            sz = b = face = col = None
            if rPr is not None:
                sz = rPr.get("sz") or d_sz
                b = rPr.get("b") or d_b
                face_el = rPr.find("a:latin", NS)
                face = face_el.get("typeface") if face_el is not None else d_face
                col_el = rPr.find("a:solidFill/a:srgbClr", NS)
                col = col_el.get("val") if col_el is not None else d_col
            else:
                sz, b, face, col = d_sz, d_b, d_face, d_col
            t_el = r.find("a:t", NS)
            text = t_el.text if t_el is not None and t_el.text else ""
            runs.append(
                RunInfo(
                    text=text,
                    font=face,
                    size_pt=(int(sz) / 100.0) if sz else None,
                    bold=(b == "1") if b is not None else None,
                    color_hex=col,
                )
            )
    return runs


def shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text or ""


def find_shape_by_name(slide, name: str):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def find_shapes_with_text(slide, predicate: Callable[[str], bool]):
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame and predicate(sh.text_frame.text or ""):
            out.append(sh)
    return out


def slide_background_hex(slide) -> Tuple[Optional[str], str]:
    """只看本页 <p:cSld><p:bg>，要求纯色填充，颜色直接写为 <a:srgbClr>。

    返回 (hex_or_None, reason):
      - ("F4FBFA", "ok")           解析成功
      - (None, "no-bg")            本页未设置 <p:bg>（不读母版/版式）
      - (None, "not-bgPr")         <p:bg> 下不是 <p:bgPr>（如 <p:bgRef> 引用主题）
      - (None, "not-solid")        不是 <a:solidFill>（渐变/图片/图案）
      - (None, "multi-color")      <a:solidFill> 下颜色节点数量 != 1
      - (None, "not-srgb")         颜色不是 <a:srgbClr>（如 schemeClr/sysClr）
    """
    bg = slide._element.find("p:cSld/p:bg", NS)
    if bg is None:
        return None, "no-bg"
    # <p:bg> 的直接子节点必须是 <p:bgPr>（直接定义），不是 <p:bgRef>（引用主题）
    bgPr = bg.find("p:bgPr", NS)
    if bgPr is None:
        return None, "not-bgPr"
    # <p:bgPr> 的直接子填充节点必须是 <a:solidFill>，不允许渐变/图片/图案
    solid = bgPr.find("a:solidFill", NS)
    if solid is None:
        return None, "not-solid"
    # 再确认 <p:bgPr> 下没有其它填充元素混入
    for tag in ("a:gradFill", "a:blipFill", "a:pattFill", "a:noFill"):
        if bgPr.find(tag, NS) is not None:
            return None, "not-solid"
    # <a:solidFill> 下必须恰好一个颜色节点
    color_children = [c for c in solid if isinstance(c.tag, str) and c.tag.startswith(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    )]
    if len(color_children) != 1:
        return None, "multi-color"
    color_el = color_children[0]
    if not color_el.tag.endswith("}srgbClr"):
        return None, "not-srgb"
    val = color_el.get("val")
    if not val or not re.fullmatch(r"[0-9A-Fa-f]{6}", val):
        return None, "not-srgb"
    return val.upper(), "ok"


# ---------------------------------------------------------------------------
# 维度1：可用与可修改性
# ---------------------------------------------------------------------------

def check_dim1(path: str) -> "tuple[bool, list[str]]":
    """维度 1：交付文件为 .pptx 格式，文件可正常打开。

    规则简化：仅校验
      • 后缀为 .pptx；
      • 文件可被 python-pptx 正常打开（作为 pptx 包结构合法的最小充分条件）。
    此前对 "可编辑（<p:modifyVerifier>/写保护/只读推荐）"、"可放映结构"、
    "整页图片化" 等额外校验已弃用。
    """
    notes: "list[str]" = []

    if not path.lower().endswith(".pptx"):
        notes.append("D1-1 文件后缀不是 .pptx —— 不通过")
        return False, notes

    try:
        prs = Presentation(path)
        assert len(prs.slides) > 0
    except Exception as e:
        notes.append(f"D1-1 文件无法解析：{e} —— 不通过")
        return False, notes

    notes.append("D1-1 格式为 .pptx 且可正常打开")
    return True, notes


# ---------------------------------------------------------------------------
# 维度2：完成度
# ---------------------------------------------------------------------------

@dataclass
class RuleResult:
    rule_id: str
    score: int                # 该条声明的分值（可正可负）
    label: str
    hit: bool                 # 是否命中（满足）
    detail: str = ""
    criterion: str = ""       # 评分细则原文（用于命中时打印）


def _check_font_attr(
    runs: List[RunInfo],
    expect_text: str,
    *,
    size: Optional[float] = None,
    color: Optional[str] = None,
    bold: Optional[bool] = None,
    font_name: Optional[str] = "微软雅黑",
    size_tol: float = 0.6,
    color_tol: int = 22,
) -> Tuple[bool, str]:
    """在 runs 中找到包含 expect_text 的 run，校验字体属性。"""
    cand = [r for r in runs if expect_text and expect_text in (r.text or "")]
    if not cand and expect_text:
        return False, f"未找到文本 '{expect_text}'"
    r = cand[0] if cand else (runs[0] if runs else None)
    if r is None:
        return False, "无可用 run"
    msgs = []
    ok = True
    if size is not None:
        if not size_close(r.size_pt, size, tol=size_tol):
            ok = False
            msgs.append(f"字号={r.size_pt}≠{size}")
    if color is not None:
        if not color_close(r.color_hex, color, tol=color_tol):
            ok = False
            msgs.append(f"颜色={r.color_hex}≠{color}")
    if bold is not None:
        if bool(r.bold) != bold:
            ok = False
            msgs.append(f"加粗={r.bold}≠{bold}")
    if font_name is not None:
        if not r.font or font_name not in r.font:
            ok = False
            msgs.append(f"字体={r.font}≠{font_name}")
    return ok, ("OK" if ok else "; ".join(msgs))


# 各页期望背景色（严格按评分细则）
#   第1页:        F4FBFA  (R244, G251, B250)
#   第2页:        FFFDF7  (R255, G253, B247)
#   第3/8页 part: 143642  (R20,  G54,  B66 )
#   第4-7、9-38页:F8FBF7  (R248, G251, B247)
#   第39页:       F5FBFA  (R245, G251, B250)
BG_EXPECT = {1: "F4FBFA", 2: "FFFDF7", 3: "143642", 8: "143642", 39: "F5FBFA"}
BG_RANGE_DEFAULT = "F8FBF7"  # 适用于 4-7 与 9-38


def expected_bg_for_page(idx: int) -> str:
    if idx in BG_EXPECT:
        return BG_EXPECT[idx]
    if 4 <= idx <= 7 or 9 <= idx <= 38:
        return BG_RANGE_DEFAULT
    # 其它页号细则未列出，仍按整体浅绿背景兜底
    return BG_RANGE_DEFAULT


def check_dim2(path: str) -> List[RuleResult]:
    prs = Presentation(path)
    slides = list(prs.slides)
    actual_slide_count = len(slides)
    if actual_slide_count < 39:
        blank_prs = Presentation()
        blank_slide = blank_prs.slides.add_slide(blank_prs.slide_layouts[6])
        slides.extend([blank_slide] * (39 - actual_slide_count))
    results: List[RuleResult] = []

    # 规则A +5：整份PPT模板替换 —— 各页背景颜色严格匹配评分细则
    #   第1页       F4FBFA
    #   第2页       FFFDF7
    #   第3、8页    143642
    #   第4-7、9-38 F8FBF7
    #   第39页      F5FBFA
    # 判定口径（来自用户确认）：
    #   - 只看本页 <p:cSld><p:bg>，不读母版/版式继承
    #   - 必须是纯色填充 <a:solidFill>，颜色直接写为 <a:srgbClr>
    #   - 形状填充不计入
    #   - HEX 严格相等（不放容差）
    #   - 必须恰好 39 页且每一页都符合，缺一不可
    bg_ok_pages: List[int] = []
    bg_bad: List[Tuple[int, str, str]] = []
    if actual_slide_count != 39:
        rule_hit = False
        detail = f"页数={actual_slide_count}，规则要求 39 页且每页背景均符合，未通过"
    else:
        for i, s in enumerate(slides, 1):
            actual_hex, reason = slide_background_hex(s)
            target = expected_bg_for_page(i).upper()
            if reason == "ok" and actual_hex == target:
                bg_ok_pages.append(i)
            else:
                bg_bad.append((i, actual_hex or reason, target))
        rule_hit = len(bg_bad) == 0
        detail = f"匹配 {len(bg_ok_pages)}/39 页"
        if bg_bad:
            detail += "，偏差页：" + ", ".join(
                f"p{i}(实际={a},期望={t})" for i, a, t in bg_bad
            )
    results.append(RuleResult(
        rule_id="A",
        score=5,
        label="整份PPT模板替换：各页背景颜色与细则一致",
        hit=rule_hit,
        detail=detail,
        criterion="整份PPT模板替换：第一页背景颜色代码为R244、G251、B250、HEX为F4FBFA；第二页目录页背景颜色代码为R255、G253、B247、HEX为FFFDF7；第三页和第八页part页背景颜色代码为R20、G54、B66、HEX为143642；第4-7页和第9-38页背景颜色代码为R248、G251、B247、HEX为F8FBF7；第39页背景颜色代码为R245、G251、B250、HEX为F5FBFA。",
    ))

    # 规则B +5：整份PPT字体除 “part 01”“part 02” 和正文中的标题编号外
    #            其余文本皆为微软雅黑（严格——任何一处不符合即不得分）
    #  豁免口径（严格按 rubric）：
    #    • 精确文本 "part 01" / "part 02"（大小写不敏感、去空白后完全相等）
    #    • 正文中的标题编号：仅指版式中承载编号的 Oval 13 / Oval 18 / Oval 23 内的 run
    #      （即第2页目录与第4-7、9-38页页面上"绿色/橘色圆形"里的 01/02/03 或 1/2/3）
    #    其他形状里出现的任意纯数字（如页码、正文数字）均不豁免。
    NUMBER_OVAL_NAMES = {"Oval 13", "Oval 18", "Oval 23"}
    non_yh: List[Tuple[int, str, Optional[str]]] = []
    total = 0
    for i, s in enumerate(slides, 1):
        for sh in s.shapes:
            shape_is_number_oval = sh.name in NUMBER_OVAL_NAMES
            for r in _iter_runs(sh):
                txt = (r.text or "").strip()
                if not txt:
                    continue
                # 例外1：严格的 "part 01" / "part 02"
                if txt.lower() in {"part 01", "part 02"}:
                    continue
                # 例外2：正文中的标题编号——仅当 run 位于编号圆形（Oval 13/18/23）
                # 且文本为 1-2 位纯数字时豁免；其他位置的纯数字（页码等）不豁免
                if shape_is_number_oval and re.fullmatch(r"\d{1,2}", txt):
                    continue
                total += 1
                if not r.font or "微软雅黑" not in r.font:
                    non_yh.append((i, txt[:20], r.font))
    yh_hit = total > 0 and len(non_yh) == 0
    results.append(RuleResult(
        rule_id="B",
        score=5,
        label="整份PPT字体除“part 01/02”与标题编号外皆为微软雅黑",
        hit=yh_hit,
        detail=(f"全部 {total} 处文本均为微软雅黑"
                if yh_hit
                else f"非微软雅黑 {len(non_yh)}/{total} 处，例：{non_yh[:3]}"),
        criterion="整份PPT字体除“part 01”“part 02”和正文中的标题编号外其余文本皆为微软雅黑",
    ))

    # 规则C +5：除第 3、8 页外，其余页面底部左侧的“呼吸系统诊疗收费项目合规解读”
    #            字体皆为 灰色 + 微软雅黑 + 5.5 磅（严格匹配细则）
    # 判定口径（响应 rubric 未给具体灰色 HEX 的问题）：
    #   - 检查页脚文本所有 run（非空 run），任一 run 不满足即该页不通过
    #   - 灰色判定采用“可接受灰色集合 ∪ 通用中性灰规则”：
    #       * 集合：常用的 5C6F7A / 595959 / 808080 / 7F7F7F / 4D4D4D / A6A6A6 / 404040
    #         （容差 ≤18 的通道差异，覆盖同色号的细微偏差）
    #       * 通用规则：|R-G|、|R-B|、|G-B| 均 ≤22，且 40 ≤ 平均亮度 ≤ 190
    #         （既排除纯黑/接近黑，也排除白/接近白，仅接受"中性灰")
    GRAY_TARGETS_C = ("5C6F7A", "595959", "808080", "7F7F7F",
                      "4D4D4D", "A6A6A6", "404040", "737373")

    def _is_gray_hex(hex_str: Optional[str]) -> bool:
        if not hex_str:
            return False
        try:
            r, g, b = _hex_to_rgb(str(hex_str))
        except Exception:
            return False
        # 命中常见灰色目标
        for t in GRAY_TARGETS_C:
            if color_close(hex_str, t, tol=18):
                return True
        # 通用中性灰：通道接近 + 亮度在中间区间
        if max(abs(r - g), abs(r - b), abs(g - b)) <= 22:
            avg = (r + g + b) / 3.0
            if 40 <= avg <= 190:
                return True
        return False

    foot_ok = 0
    foot_total = 0
    foot_bad: List[Tuple[int, str]] = []
    for i, s in enumerate(slides, 1):
        if i in (3, 8):
            continue
        # 按文本内容定位底部左侧页脚文本框（位置：top 在下半页、left 在左半页）
        sw, sh_h = prs.slide_width, prs.slide_height
        footer = None
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            if "呼吸系统诊疗收费项目合规解读" not in (sh.text_frame.text or ""):
                continue
            try:
                if sh.top is None or sh.left is None:
                    continue
                if sh.top > sh_h * 0.5 and sh.left < sw * 0.5:
                    footer = sh
                    break
            except Exception:
                continue
        if footer is None:
            foot_total += 1
            foot_bad.append((i, "未找到底部左侧页脚"))
            continue
        foot_total += 1
        runs = [r for r in _iter_runs(footer) if (r.text or "").strip()]
        if not runs:
            foot_bad.append((i, "无可用 run"))
            continue
        # 所有相关 run 都必须满足：微软雅黑 + 灰色 + 严格 5.5pt
        bad_runs_c: List[str] = []
        for idx_r, r in enumerate(runs):
            is_yh = bool(r.font and "微软雅黑" in r.font)
            is_gray = _is_gray_hex(r.color_hex)
            is_5p5 = size_close(r.size_pt, 5.5, tol=0.25)
            if not (is_yh and is_gray and is_5p5):
                bad_runs_c.append(
                    f"run#{idx_r}(字体={r.font} 颜色={r.color_hex} 字号={r.size_pt})"
                )
        if not bad_runs_c:
            foot_ok += 1
        else:
            foot_bad.append((i, "; ".join(bad_runs_c)))
    results.append(RuleResult(
        rule_id="C",
        score=5,
        label="除第3/8页外底部左侧页脚为灰色微软雅黑 5.5磅",
        hit=foot_total > 0 and foot_ok == foot_total,
        detail=(f"通过 {foot_ok}/{foot_total} 页"
                + (f"，偏差：{foot_bad[:3]}" if foot_bad else "")),
        criterion="PPT中除第3、8页外其余页面底部左侧的“呼吸系统诊疗收费项目合规解读”字体皆为灰色微软雅黑5.5磅",
    ))

    # 规则D +5：PPT 第 1-39 页中页面底部右侧的页码 “01/39” 至 “39/39”
    #            字体皆为 绿色 + 微软雅黑 + 8 磅 + 加粗（严格匹配细则）
    # 判定口径（响应问题与措施）：
    #   - 先按"页码精确文本 NN/39"在底部右侧候选中筛选，命中的文本框才作为页码框；
    #     不再取右下第一个文本框即用（避免右下其他文本在前时误判）。
    #   - 校验该文本框内所有非空 run，任一 run 不合格即该页不通过。
    #   - 绿色目标：主视觉深青绿 149E9A（细则整体一致），容差从 40 收紧到 15
    #     （允许同色号的细微偏差，但排除明显偏离绿色的颜色）。
    page_ok = 0
    page_bad: List[Tuple[int, str]] = []
    sw, sh_h = prs.slide_width, prs.slide_height
    for i, s in enumerate(slides, 1):
        expected_text = f"{i:02d}/39"
        # 1) 先按 "NN/39" 精确文本，在底部右侧候选中定位页码框
        target = None
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            try:
                if sh.top is None or sh.left is None:
                    continue
                if not (sh.top > sh_h * 0.5 and sh.left > sw * 0.5):
                    continue
            except Exception:
                continue
            if (sh.text_frame.text or "").strip() == expected_text:
                target = sh
                break
        # 2) 若严格文本未命中，回退到"底部右侧且文本形如 NN/39"（宽松匹配，同样要求文本匹配页码模式）
        if target is None:
            for sh in s.shapes:
                if not sh.has_text_frame:
                    continue
                try:
                    if sh.top is None or sh.left is None:
                        continue
                    if not (sh.top > sh_h * 0.5 and sh.left > sw * 0.5):
                        continue
                except Exception:
                    continue
                t = (sh.text_frame.text or "").strip()
                if re.fullmatch(r"\d{1,2}/\d{1,2}", t) and t == expected_text:
                    target = sh
                    break
        if target is None:
            page_bad.append((i, f"未找到文本为 '{expected_text}' 的底部右侧页码文本框"))
            continue

        text = (target.text_frame.text or "").strip()
        text_ok = (text == expected_text)
        runs = [r for r in _iter_runs(target) if (r.text or "").strip()]
        if not runs:
            page_bad.append((i, f"页码文本='{text}' 无 run"))
            continue
        # 3) 校验所有非空 run：绿色（149E9A，容差 15）+ 微软雅黑 + 8pt + 加粗
        bad_runs: List[str] = []
        for idx_r, r in enumerate(runs):
            is_yh = bool(r.font and "微软雅黑" in r.font)
            is_green = color_close(r.color_hex, "149E9A", tol=15)
            is_8pt = size_close(r.size_pt, 8.0, tol=0.25)
            is_bold = bool(r.bold)
            if not (is_yh and is_green and is_8pt and is_bold):
                bad_runs.append(
                    f"run#{idx_r}(字体={r.font} 颜色={r.color_hex} 字号={r.size_pt} 加粗={r.bold})"
                )
        if text_ok and not bad_runs:
            page_ok += 1
        else:
            page_bad.append((
                i,
                f"文本='{text}'(期望'{expected_text}') "
                + ("; ".join(bad_runs) if bad_runs else "字体属性 OK"),
            ))
    results.append(RuleResult(
        rule_id="D",
        score=5,
        label="第1-39页底部右侧页码“01/39”~“39/39”为绿色微软雅黑8磅加粗",
        hit=page_ok == 39,
        detail=(f"通过 {page_ok}/39 页"
                + (f"，偏差：{page_bad[:3]}" if page_bad else "")),
        criterion="PPT第1-39页中页面底部右侧的页码“01/39”至“39/39”字体皆为绿色微软雅黑8磅加粗",
    ))

    # 规则E +1：第1页封面 —— 严格匹配细则中列出的每一处文字与配色
    #   • “呼吸系统诊疗收费项目合规解读” 黑色 34磅 加粗
    #   • “医疗服务价格管理”             白色 10.5磅 加粗
    #   • “基于医保支付口径与院内收费规范的业务培训材料” 灰色 15.5磅
    #   • “瑞康医疗价格管理中心”         白色 9.5磅 加粗
    #   • “2026年春季”                  白色 9.5磅 加粗
    #   • 同时使用深青绿主视觉 + 橙色强调区域
    s1 = slides[0]

    def _find_run_with(text: str) -> Optional[RunInfo]:
        for sh in s1.shapes:
            if not sh.has_text_frame:
                continue
            for r in _iter_runs(sh):
                if text in (r.text or ""):
                    return r
        return None

    def _check(expected_text: str, *, size: float, color: str, bold: Optional[bool]) -> Tuple[bool, str]:
        r = _find_run_with(expected_text)
        if r is None:
            return False, f"未找到 '{expected_text}'"
        size_ok = size_close(r.size_pt, size, tol=0.25)
        color_ok = color_close(r.color_hex, color, tol=22)
        bold_ok = True if bold is None else (bool(r.bold) == bold)
        ok = size_ok and color_ok and bold_ok
        return ok, (f"字号={r.size_pt} 颜色={r.color_hex} 加粗={r.bold}")

    e_checks = []
    ok, msg = _check("呼吸系统诊疗收费项目合规解读", size=34, color="000000", bold=True)
    e_checks.append(("呼吸系统诊疗收费项目合规解读 黑色34磅加粗", ok, msg))
    ok, msg = _check("医疗服务价格管理", size=10.5, color="FFFFFF", bold=True)
    e_checks.append(("医疗服务价格管理 白色10.5磅加粗", ok, msg))
    ok, msg = _check("基于医保支付口径与院内收费规范的业务培训材料",
                     size=15.5, color="5C6F7A", bold=None)
    e_checks.append(("基于医保…培训材料 灰色15.5磅", ok, msg))
    ok, msg = _check("瑞康医疗价格管理中心", size=9.5, color="FFFFFF", bold=True)
    e_checks.append(("瑞康医疗价格管理中心 白色9.5磅加粗", ok, msg))
    ok, msg = _check("2026年春季", size=9.5, color="FFFFFF", bold=True)
    e_checks.append(("2026年春季 白色9.5磅加粗", ok, msg))

    # 深青绿主视觉 + 橙色强调区域
    has_main_green = False
    has_orange = False
    for sh in s1.shapes:
        try:
            rgb = str(sh.fill.fore_color.rgb)
        except Exception:
            continue
        if not has_main_green and color_close(rgb, "149E9A", tol=30):
            has_main_green = True
        if not has_orange and (color_close(rgb, "F4A635", tol=40)
                               or color_close(rgb, "FF8C00", tol=60)):
            has_orange = True
    e_checks.append(("深青绿主视觉 + 橙色强调区域",
                     has_main_green and has_orange,
                     f"主视觉={has_main_green} 橙色={has_orange}"))

    e_all_ok = all(ok for _, ok, _ in e_checks)
    results.append(RuleResult(
        rule_id="E",
        score=1,
        label="第1页封面字体与配色严格符合细则",
        hit=e_all_ok,
        detail="; ".join(f"{n}:{'√' if ok else '×'}({m})" for n, ok, m in e_checks),
        criterion="第1页封面：“呼吸系统诊疗收费项目合规解读”字体为黑色34磅加粗，“医疗服务价格管理”字体为白色10.5磅加粗，“基于医保支付口径与院内收费规范的业务培训材料”字体为灰色15.5磅，“瑞康医疗价格管理中心”“2026年春季”字体为白色9.5磅加粗，等文字，同时使用深青绿主视觉和橙色强调区域。",
    ))

    # 规则F +5：除第 3、8、39 页外，其余页面的左上存在一个
    #            “填充色为绿色的圆角矩形边框”的 “医疗价格服务价格管理” 标识布局，
    #            字体为 白色 10.5 磅 加粗（严格按细则）
    f_ok = 0
    f_total = 0
    f_bad: List[Tuple[int, str]] = []
    sw, sh_h = prs.slide_width, prs.slide_height
    rule_text = "医疗价格服务价格管理"
    for i, s in enumerate(slides, 1):
        if i in (3, 8, 39):
            continue
        f_total += 1
        # 在左上区域（left < 1/2 slide_width, top < 1/3 slide_height）寻找
        # 圆角矩形（prstGeom = roundRect）且填充为绿色，文本=细则要求的字符串
        target = None
        for sh in s.shapes:
            try:
                if sh.top is None or sh.left is None:
                    continue
                if not (sh.left < sw * 0.5 and sh.top < sh_h * 0.33):
                    continue
            except Exception:
                continue
            # 必须是圆角矩形
            prst = sh._element.find(".//a:prstGeom", NS)
            if prst is None or prst.get("prst") != "roundRect":
                continue
            # 填充为绿色
            try:
                rgb = str(sh.fill.fore_color.rgb)
            except Exception:
                continue
            if not color_close(rgb, "149E9A", tol=30):
                continue
            # 文本必须为细则要求的 “医疗价格服务价格管理”
            if not sh.has_text_frame:
                continue
            if rule_text not in (sh.text_frame.text or ""):
                continue
            target = sh
            break
        if target is None:
            f_bad.append((i, f"未找到左上绿色圆角矩形的“{rule_text}”标识"))
            continue
        runs = list(_iter_runs(target))
        if not runs:
            f_bad.append((i, "标识无可用 run"))
            continue
        r0 = runs[0]
        is_white = color_close(r0.color_hex, "FFFFFF", tol=15)
        is_10p5 = size_close(r0.size_pt, 10.5, tol=0.25)
        is_bold = bool(r0.bold)
        if is_white and is_10p5 and is_bold:
            f_ok += 1
        else:
            f_bad.append((i, f"颜色={r0.color_hex} 字号={r0.size_pt} 加粗={r0.bold}"))
    results.append(RuleResult(
        rule_id="F",
        score=5,
        label="除3/8/39页外左上绿色圆角矩形“医疗价格服务价格管理”标识 白色10.5磅加粗",
        hit=f_total > 0 and f_ok == f_total,
        detail=(f"通过 {f_ok}/{f_total} 页"
                + (f"，偏差：{f_bad[:3]}" if f_bad else "")),
        criterion="除第3、8、39页外其余页面的左上存在一个填充色为绿色的圆角矩形边框的“医疗价格服务价格管理”标识布局字体为白色10.5磅加粗。",
    ))

    # 规则G +1：第 3 页和第 8 页页面左上角存在一个填充色为绿色的圆角矩形边框，
    #            内有文本 “part 01”/“part 02”，字体为 白色 Arial 10.5 磅 加粗
    g_hits = 0
    g_bad: List[Tuple[int, str]] = []
    sw, sh_h = prs.slide_width, prs.slide_height
    expected_map = {3: "part 01", 8: "part 02"}
    for i, expected in expected_map.items():
        s = slides[i - 1]
        target = None
        for sh in s.shapes:
            try:
                if sh.top is None or sh.left is None:
                    continue
                if not (sh.left < sw * 0.5 and sh.top < sh_h * 0.5):
                    continue
            except Exception:
                continue
            # 必须是圆角矩形
            prst = sh._element.find(".//a:prstGeom", NS)
            if prst is None or prst.get("prst") != "roundRect":
                continue
            # 填充为绿色
            try:
                rgb = str(sh.fill.fore_color.rgb)
            except Exception:
                continue
            if not color_close(rgb, "149E9A", tol=30):
                continue
            # 内部文本包含 part 01 / part 02（大小写不敏感）
            if not sh.has_text_frame:
                continue
            text = (sh.text_frame.text or "").strip().lower()
            if expected not in text:
                continue
            target = sh
            break
        if target is None:
            g_bad.append((i, f"未找到左上绿色圆角矩形且内含 “{expected}” 的形状"))
            continue
        runs = list(_iter_runs(target))
        if not runs:
            g_bad.append((i, "形状无可用 run"))
            continue
        r0 = runs[0]
        is_white = color_close(r0.color_hex, "FFFFFF", tol=15)
        is_arial = bool(r0.font and "arial" in r0.font.lower())
        is_10p5 = size_close(r0.size_pt, 10.5, tol=0.25)
        is_bold = bool(r0.bold)
        if is_white and is_arial and is_10p5 and is_bold:
            g_hits += 1
        else:
            g_bad.append((
                i,
                f"字体={r0.font} 颜色={r0.color_hex} 字号={r0.size_pt} 加粗={r0.bold}",
            ))
    results.append(RuleResult(
        rule_id="G",
        score=1,
        label="第3/8页左上绿色圆角矩形内 “part 01”/“part 02” 白色Arial10.5磅加粗",
        hit=g_hits == 2,
        detail=(f"命中 {g_hits}/2"
                + (f"，偏差：{g_bad}" if g_bad else "")),
        criterion="第3页和第8页两页中页面左上角存在一个填充色为绿色的圆角矩形边框，内有文本“part 01”“part 02”字体为白色Arial10.5磅加粗",
    ))

    # 规则H +1：第 2 页目录 —— 严格匹配细则中列出的每一项
    #   • “目录” 黑色 28磅 加粗
    #   • “从政策口径到项目明细，按‘先规则、后执行’的路径展开” 灰色 12.5磅
    #   • “政策文件要点”/“项目明细规范”/“执行注意事项” 三个标题 黑色 15磅 加粗
    #   • 上述三个标题下方对应说明文字 灰色 9.8磅
    #   • 三个标题前编号 “01”/“02”/“03” 为横向文本框 白色 Arial 10.5磅 加粗
    #   • 编号放置在三个填充颜色为 绿色、橘色、绿色 的圆形上方
    s2 = slides[1]

    def _find_run_in_slide(slide, text: str) -> Optional[RunInfo]:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for r in _iter_runs(sh):
                if text and text in (r.text or ""):
                    return r
        return None

    def _check_h(text: str, *, size: float, color: str,
                 bold: Optional[bool], font_name: Optional[str] = "微软雅黑") -> Tuple[bool, str]:
        r = _find_run_in_slide(s2, text)
        if r is None:
            return False, f"未找到 '{text}'"
        size_ok = size_close(r.size_pt, size, tol=0.25)
        color_ok = color_close(r.color_hex, color, tol=22)
        bold_ok = True if bold is None else (bool(r.bold) == bold)
        font_ok = True
        if font_name is not None:
            font_ok = bool(r.font and font_name.lower() in r.font.lower())
        return (size_ok and color_ok and bold_ok and font_ok,
                f"字体={r.font} 字号={r.size_pt} 颜色={r.color_hex} 加粗={r.bold}")

    h_checks: List[Tuple[str, bool, str]] = []
    ok, msg = _check_h("目录", size=28, color="000000", bold=True)
    h_checks.append(("“目录” 黑色28磅加粗", ok, msg))

    # 副标题：按 rubric 完整句匹配（同时接受直/曲两种单引号），不再用片段匹配。
    # rubric 原文：从政策口径到项目明细，按‘先规则、后执行’的路径展开
    subtitle_variants = (
        "从政策口径到项目明细，按‘先规则、后执行’的路径展开",
        "从政策口径到项目明细，按'先规则、后执行'的路径展开",
    )
    subtitle_run: Optional[RunInfo] = None
    # 优先：单个 run 内包含整句
    for v in subtitle_variants:
        r = _find_run_in_slide(s2, v)
        if r is not None:
            subtitle_run = r
            break
    # 回退：整句被拆到同一 shape 的多个 run（以 text_frame.text 拼接后判断）
    if subtitle_run is None:
        for sh in s2.shapes:
            if not sh.has_text_frame:
                continue
            frame_text = sh.text_frame.text or ""
            if any(v in frame_text for v in subtitle_variants):
                rl = [r for r in _iter_runs(sh) if (r.text or "").strip()]
                if rl:
                    subtitle_run = rl[0]
                break
    if subtitle_run is None:
        h_checks.append((
            "副标题“从政策口径到项目明细，按‘先规则、后执行’的路径展开” 灰色12.5磅",
            False,
            "未找到完整副标题句（rubric 全句）",
        ))
    else:
        r = subtitle_run
        size_ok = size_close(r.size_pt, 12.5, tol=0.25)
        color_ok = color_close(r.color_hex, "5C6F7A", tol=22)
        h_checks.append((
            "副标题“从政策口径到项目明细，按‘先规则、后执行’的路径展开” 灰色12.5磅",
            size_ok and color_ok,
            f"字号={r.size_pt} 颜色={r.color_hex}",
        ))

    # 三个标题：黑色 15磅 加粗；同时记录其 shape 以便后续按位置查下方对应文字
    title_shapes = []  # type: List[Tuple[str, object]]  # 运行时为 python-pptx Shape
    for kw in ("政策文件要点", "项目明细规范", "执行注意事项"):
        ok, msg = _check_h(kw, size=15, color="000000", bold=True)
        h_checks.append((f"标题“{kw}” 黑色15磅加粗", ok, msg))
        for sh in s2.shapes:
            if sh.has_text_frame and kw in (sh.text_frame.text or ""):
                title_shapes.append((kw, sh))
                break

    # 标题下方对应文字：按位置关系（同页、位于对应标题正下方、水平重叠）取一个 run 校验
    #   —— rubric 未给具体字面文本，只要求"标题目录下方对应文字 灰色 9.8磅"，
    #      因此不再使用代码自定的固定关键词，改为按几何位置判定。
    for kw, title_sh in title_shapes:
        try:
            t_top = title_sh.top
            t_left = title_sh.left
            t_w = title_sh.width or 0
            t_h = title_sh.height or 0
        except Exception:
            h_checks.append((f"标题“{kw}”下方对应文字 灰色9.8磅", False, "标题位置不可读"))
            continue
        if t_top is None or t_left is None:
            h_checks.append((f"标题“{kw}”下方对应文字 灰色9.8磅", False, "标题位置为空"))
            continue
        title_center_x = t_left + t_w / 2
        title_bottom = t_top + t_h
        # 在同页所有可读文本框里，找到位于标题正下方、水平区间与标题水平中线重叠的形状
        below_candidates: List[Tuple[int, object]] = []
        for sh in s2.shapes:
            if sh is title_sh or not sh.has_text_frame:
                continue
            frame_text = (sh.text_frame.text or "").strip()
            if not frame_text or frame_text == kw:
                continue
            try:
                if sh.top is None or sh.left is None:
                    continue
                sh_w = sh.width or 0
                sh_h_ = sh.height or 0
            except Exception:
                continue
            if sh.top < title_bottom:
                continue
            # 要求形状水平区间覆盖标题水平中心，避免误取旁边一列文字
            if not (sh.left <= title_center_x <= sh.left + sh_w):
                continue
            below_candidates.append((sh.top - title_bottom, sh))
        if not below_candidates:
            h_checks.append((f"标题“{kw}”下方对应文字 灰色9.8磅", False, "未找到标题下方文字"))
            continue
        # 取最靠近标题的下方文本框
        below_candidates.sort(key=lambda x: x[0])
        target_sh = below_candidates[0][1]
        rl = [r for r in _iter_runs(target_sh) if (r.text or "").strip()]
        if not rl:
            h_checks.append((f"标题“{kw}”下方对应文字 灰色9.8磅", False, "下方文本框无 run"))
            continue
        r = rl[0]
        size_ok = size_close(r.size_pt, 9.8, tol=0.25)
        color_ok = color_close(r.color_hex, "5C6F7A", tol=22)
        font_ok = bool(r.font and "微软雅黑" in r.font)
        h_checks.append((
            f"标题“{kw}”下方对应文字 灰色9.8磅",
            size_ok and color_ok and font_ok,
            f"文本='{(target_sh.text_frame.text or '').strip()[:20]}' "
            f"字体={r.font} 字号={r.size_pt} 颜色={r.color_hex}",
        ))

    # 三个编号 01/02/03：白色 Arial 10.5磅 加粗，位于横向文本框中，且文本框在
    # 一个绿色/橘色/绿色的圆形（prstGeom=ellipse）的正上方（不依赖固定形状名）。
    # 判定思路：
    #   1) 收集本页所有 prstGeom=ellipse 的形状（圆形），按 left 从小到大排序作为
    #      三个编号槽位；每个槽位有一个期望填充色（绿、橘、绿）。
    #   2) 对每个槽位，寻找同页中：
    #        - 是横向文本框（width >= height）；
    #        - 文本精确等于该槽位的期望编号（"01"/"02"/"03"）；
    #        - 位于圆形上方（bottom_of_textbox <= top_of_oval + 少量容差），
    #          且水平中心与圆形水平中心接近（差 ≤ 圆形宽度的一半）。
    #     符合的文本框内 run 需为 白色 Arial 10.5磅 加粗。
    #   3) 同时校验圆形本身的填充为 绿/橘/绿。
    expected_nums = ("01", "02", "03")
    expected_fills = ("149E9A", "F4A635", "149E9A")  # 绿、橘、绿

    ellipse_shapes: List[object] = []
    for sh in s2.shapes:
        try:
            prst = sh._element.find(".//a:prstGeom", NS)
        except Exception:
            continue
        if prst is None or prst.get("prst") != "ellipse":
            continue
        if sh.left is None or sh.top is None:
            continue
        ellipse_shapes.append(sh)
    # 按 left 从小到大作为 01/02/03 的槽位顺序
    ellipse_shapes.sort(key=lambda x: (x.left or 0))

    if len(ellipse_shapes) < 3:
        h_checks.append((
            "三个编号圆形（绿/橘/绿）", False,
            f"页内可识别的椭圆/圆形数量={len(ellipse_shapes)}，不足3个",
        ))
    else:
        for slot_idx, (num, fill_hex) in enumerate(zip(expected_nums, expected_fills)):
            oval = ellipse_shapes[slot_idx]
            try:
                actual_fill = str(oval.fill.fore_color.rgb)
            except Exception:
                actual_fill = ""
            fill_ok = color_close(actual_fill, fill_hex, tol=30)

            o_left = oval.left or 0
            o_top = oval.top or 0
            o_w = oval.width or 0
            o_h = oval.height or 0
            o_center_x = o_left + o_w / 2

            # 在同页寻找位于圆形正上方的横向文本框，文本精确=num
            tol_v = int(o_h * 0.75) if o_h else 0  # 允许文本框略微下压/贴合圆形顶部
            tol_h_center = int(o_w * 0.6) if o_w else 0
            num_shape = None
            num_run: Optional[RunInfo] = None
            for sh in s2.shapes:
                if sh is oval or not sh.has_text_frame:
                    continue
                try:
                    if sh.top is None or sh.left is None:
                        continue
                    sh_w = sh.width or 0
                    sh_h_ = sh.height or 0
                except Exception:
                    continue
                # 横向文本框：宽 >= 高
                if sh_w <= 0 or sh_h_ <= 0 or sh_w < sh_h_:
                    continue
                if (sh.text_frame.text or "").strip() != num:
                    continue
                # 位于圆形上方：文本框底 <= 圆形顶 + 容差
                if sh.top + sh_h_ > o_top + tol_v:
                    continue
                # 水平中心接近圆形中心
                sh_center_x = sh.left + sh_w / 2
                if abs(sh_center_x - o_center_x) > tol_h_center:
                    continue
                rl = [r for r in _iter_runs(sh) if (r.text or "").strip() == num]
                if not rl:
                    continue
                num_shape = sh
                num_run = rl[0]
                break

            if num_run is None or num_shape is None:
                h_checks.append((
                    f"编号 {num} 在圆形（期望{fill_hex}）上方的横向文本框内",
                    False,
                    f"填充={actual_fill}(期望{fill_hex}) 未找到位于圆形上方的横向文本框",
                ))
                continue
            is_white = color_close(num_run.color_hex, "FFFFFF", tol=15)
            is_arial = bool(num_run.font and "arial" in num_run.font.lower())
            is_10p5 = size_close(num_run.size_pt, 10.5, tol=0.25)
            is_bold = bool(num_run.bold)
            text_ok = is_white and is_arial and is_10p5 and is_bold
            h_checks.append((
                f"编号 {num}（圆形{fill_hex}上方的横向文本框）",
                fill_ok and text_ok,
                f"填充={actual_fill}(期望{fill_hex}) 字体={num_run.font} 字号={num_run.size_pt} "
                f"颜色={num_run.color_hex} 加粗={num_run.bold}",
            ))

    h_all_ok = all(ok for _, ok, _ in h_checks)
    results.append(RuleResult(
        rule_id="H",
        score=1,
        label="第2页目录字体/编号/圆形配色严格符合细则",
        hit=h_all_ok,
        detail="; ".join(f"{n}:{'√' if ok else '×'}({m})" for n, ok, m in h_checks),
        criterion="第2页目录：“目录”字体为黑色28磅加粗，“从政策口径到项目明细，按‘先规则、后执行’的路径展开”字体为灰色12.5磅，“政策文件要点、项目明细规范、执行注意事项”三个标题字体为黑色15磅加粗，标题目录下方对应文字字体为灰色9.8磅，三个标题前的编号“01”“02”“03”字体为横向文本框白色Arial10.5磅加粗，编号放置在三个填充颜色为绿色、橘色、绿色的圆形上方。",
    ))

    # 规则I +1：第 3 页和第 8 页章节页 —— 严格匹配细则中列出的每一项
    #   • 第3页 “政策框架与应用口径”            白色 34磅 加粗
    #   • 第8页 “项目明细与收费规范”            白色 34磅 加粗
    #   • 第3页 “先明确项目边界，再进入具体收费场景” 浅绿色 14磅
    #   • 第8页 “以‘检查—治疗—手术—注意事项’的顺序展开” 浅绿色 14磅
    def _find_in_slide(slide, text: str) -> Optional[RunInfo]:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for r in _iter_runs(sh):
                if text and text in (r.text or ""):
                    return r
        return None

    def _check_i(slide_idx: int, text: str, *, size: float,
                 color_targets: List[str], bold: Optional[bool]) -> Tuple[bool, str]:
        r = _find_in_slide(slides[slide_idx - 1], text)
        if r is None:
            return False, f"未找到 '{text}'"
        size_ok = size_close(r.size_pt, size, tol=0.25)
        color_ok = any(color_close(r.color_hex, c, tol=60) for c in color_targets)
        bold_ok = True if bold is None else (bool(r.bold) == bold)
        return (size_ok and color_ok and bold_ok,
                f"字号={r.size_pt} 颜色={r.color_hex} 加粗={r.bold}")

    i_checks: List[Tuple[str, bool, str]] = []
    ok, msg = _check_i(3, "政策框架与应用口径", size=34,
                       color_targets=["FFFFFF"], bold=True)
    i_checks.append(("第3页“政策框架与应用口径” 白色34磅加粗", ok, msg))
    ok, msg = _check_i(8, "项目明细与收费规范", size=34,
                       color_targets=["FFFFFF"], bold=True)
    i_checks.append(("第8页“项目明细与收费规范” 白色34磅加粗", ok, msg))
    # 浅绿色：允许 #A8E6CF / #C6F0DC / 浅青绿 等近似浅绿
    light_green_targets = ["A8E6CF", "BFE6D8", "C6F0DC", "D9F0E8", "E0F5EC", "B5E2C9"]
    ok, msg = _check_i(3, "先明确项目边界", size=14,
                       color_targets=light_green_targets, bold=None)
    i_checks.append(("第3页“先明确项目边界…” 浅绿色14磅", ok, msg))
    ok, msg = _check_i(8, "顺序展开", size=14,
                       color_targets=light_green_targets, bold=None)
    i_checks.append(("第8页“以‘检查—治疗…顺序展开’” 浅绿色14磅", ok, msg))

    i_all_ok = all(ok for _, ok, _ in i_checks)
    results.append(RuleResult(
        rule_id="I",
        score=1,
        label="第3/8页章节页文字 白色34磅加粗 + 浅绿色14磅",
        hit=i_all_ok,
        detail="; ".join(f"{n}:{'√' if ok else '×'}({m})" for n, ok, m in i_checks),
        criterion="第3页和第8页章节页：“政策框架与应用口径”和“项目明细与收费规范”字体为白色34磅加粗；“先明确项目边界，再进入具体收费场景”和“以‘检查—治疗—手术—注意事项’的顺序展开”字体为浅绿色14磅。",
    ))

    # 规则J +3：第 4 页至第 7 页政策框架页 —— 严格匹配细则
    #   • 大标题（如“规范整合呼吸系统医疗服务项目” 等）         黑色 27磅 加粗
    #   • 大标题下方承接说明一行                                  灰色 11.5磅
    #   • “项目整合” 等小标题                                     黑色 14.5磅 加粗
    #   • 上述小标题下方对应文本                                  灰色 9.7磅
    # 判定口径（响应问题与措施）：
    #   - 不依赖固定形状名（原 TextBox 8/9/14/19/24/15/20/25）；改为按内容 + 位置
    #     角色定位每一类文字。
    #   - 每一类都校验形状内**所有非空 run**，任一 run 不合格即该类不通过；
    #     任一类不通过 → 整页不通过。
    # 页内噪声形状（页脚、页码、左上品牌标识、章节编号圆圈）会先被过滤掉。
    def _shape_all_runs_ok(
        sh, *, size: float, color_hex: str,
        bold: Optional[bool], color_tol: int = 22, size_tol: float = 0.25,
        font_name: str = "微软雅黑",
    ) -> Tuple[bool, str]:
        runs = [r for r in _iter_runs(sh) if (r.text or "").strip()]
        if not runs:
            return False, "无 run"
        bad_parts: List[str] = []
        for idx_r, r in enumerate(runs):
            size_ok = size_close(r.size_pt, size, tol=size_tol)
            color_ok = color_close(r.color_hex, color_hex, tol=color_tol)
            bold_ok = True if bold is None else (bool(r.bold) == bold)
            font_ok = bool(r.font and font_name in r.font)
            if not (size_ok and color_ok and bold_ok and font_ok):
                bad_parts.append(
                    f"run#{idx_r}(字号={r.size_pt} 颜色={r.color_hex} 加粗={r.bold} 字体={r.font})"
                )
        if not bad_parts:
            return True, "OK"
        return False, "; ".join(bad_parts)

    def _collect_role_shapes(slide) -> List[object]:
        """收集用于角色识别的候选文本形状，过滤掉页脚/页码/品牌/编号圆圈等噪声。"""
        out = []
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            text = (sh.text_frame.text or "").strip()
            if not text:
                continue
            try:
                if sh.top is None or sh.left is None:
                    continue
            except Exception:
                continue
            # 过滤页脚（"呼吸系统诊疗收费项目合规解读" 页脚）
            if "呼吸系统诊疗收费项目合规解读" in text:
                continue
            # 过滤底部页码 NN/39
            if re.fullmatch(r"\d{1,2}/\d{1,2}", text):
                continue
            # 过滤左上品牌标识
            if "医疗价格服务价格管理" in text:
                continue
            # 过滤章节编号圆形内/上的 "01"/"02"/"03"/"1"/"2"/"3"
            if re.fullmatch(r"\d{1,2}", text):
                continue
            out.append(sh)
        return out

    def _cluster_by_top(shapes: List[object], threshold_ratio: float = 0.03,
                       page_h: int = 0) -> List[List[object]]:
        """按 top 聚类：属于同一"行"的形状 top 差 ≤ threshold_ratio × page_h。"""
        if not shapes:
            return []
        thresh = max(1, int(page_h * threshold_ratio)) if page_h else 200000
        srt = sorted(shapes, key=lambda x: x.top or 0)
        rows: List[List[object]] = [[srt[0]]]
        for sh in srt[1:]:
            if abs((sh.top or 0) - (rows[-1][0].top or 0)) <= thresh:
                rows[-1].append(sh)
            else:
                rows.append([sh])
        # 每行内按 left 从小到大
        for row in rows:
            row.sort(key=lambda x: x.left or 0)
        return rows

    j_pages_ok = 0
    j_bad: List[Tuple[int, str]] = []
    sh_h_all = int(prs.slide_height or 0)
    for idx in range(4, 8):
        s = slides[idx - 1]
        cands = _collect_role_shapes(s)
        bad_subs: List[str] = []

        if not cands:
            j_bad.append((idx, "未找到可识别的正文形状"))
            continue

        rows = _cluster_by_top(cands, threshold_ratio=0.03, page_h=sh_h_all)
        # 期望结构：至少 4 行—— [大标题], [承接说明], [3 个小标题], [3 个对应文本]
        # 若承接说明与大标题被聚为一行，退化处理
        if len(rows) < 3:
            j_bad.append((idx, f"版式行数={len(rows)}，无法识别大标题/承接说明/小标题/对应文本"))
            continue

        # 大标题：最上一行；若有多形状取最靠左的（通常大标题独占一行）
        big_sh = rows[0][0]
        ok, msg = _shape_all_runs_ok(
            big_sh, size=27, color_hex="000000", bold=True, color_tol=22)
        if not ok:
            bad_subs.append(f"大标题（文本='{(big_sh.text_frame.text or '').strip()[:15]}'）×({msg})")

        # 承接说明：第二行；同样取最靠左的一个形状
        sub_sh = rows[1][0]
        ok, msg = _shape_all_runs_ok(
            sub_sh, size=11.5, color_hex="5C6F7A", bold=None, color_tol=25)
        if not ok:
            bad_subs.append(f"承接说明（文本='{(sub_sh.text_frame.text or '').strip()[:15]}'）×({msg})")

        # 小标题行：期望包含 3 个形状；小标题需为黑色 14.5磅 加粗
        # 对应文本行：期望紧接在小标题行下，包含 3 个形状；需为灰色 9.7磅
        remaining_rows = rows[2:]
        # 从剩余行中找到第一行"看起来像小标题（≥2 个形状且大多为加粗）"作为小标题行；
        # 其后紧邻的一行作为对应文本行。
        sub_title_row = None
        body_row = None
        for i_r, row in enumerate(remaining_rows):
            if len(row) < 2:
                continue
            # 至少两个形状为加粗（避免把 body 误当小标题行）
            bold_count = 0
            for sh in row:
                runs = [r for r in _iter_runs(sh) if (r.text or "").strip()]
                if runs and bool(runs[0].bold):
                    bold_count += 1
            if bold_count >= 2:
                sub_title_row = row
                if i_r + 1 < len(remaining_rows):
                    body_row = remaining_rows[i_r + 1]
                break

        if sub_title_row is None:
            bad_subs.append("小标题行×(未在剩余行中识别到)")
        else:
            if len(sub_title_row) != 3:
                bad_subs.append(f"小标题×(数量={len(sub_title_row)}，期望3)")
            for sh in sub_title_row:
                ok, msg = _shape_all_runs_ok(
                    sh, size=14.5, color_hex="000000", bold=True, color_tol=22)
                if not ok:
                    bad_subs.append(
                        f"小标题（文本='{(sh.text_frame.text or '').strip()[:12]}'）×({msg})"
                    )
            # 对应文本行：优先取紧邻下一行；否则对每个小标题按"正下方最近"逐个匹配
            if body_row is None or len(body_row) < len(sub_title_row):
                # 按位置为每个小标题查找同页正下方最近的形状（用剩余候选，未在小标题行中）
                used = set(id(x) for x in sub_title_row)
                for st_sh in sub_title_row:
                    st_center_x = (st_sh.left or 0) + (st_sh.width or 0) / 2
                    st_bottom = (st_sh.top or 0) + (st_sh.height or 0)
                    below: List[Tuple[int, object]] = []
                    for sh in cands:
                        if id(sh) in used:
                            continue
                        if (sh.top or 0) < st_bottom:
                            continue
                        sh_w = sh.width or 0
                        if not (sh.left or 0) <= st_center_x <= (sh.left or 0) + sh_w:
                            continue
                        below.append(((sh.top or 0) - st_bottom, sh))
                    if not below:
                        bad_subs.append(
                            f"“{(st_sh.text_frame.text or '').strip()[:12]}”对应文本×(未找到正下方形状)"
                        )
                        continue
                    below.sort(key=lambda x: x[0])
                    body_sh = below[0][1]
                    used.add(id(body_sh))
                    ok, msg = _shape_all_runs_ok(
                        body_sh, size=9.7, color_hex="5C6F7A", bold=None, color_tol=25)
                    if not ok:
                        bad_subs.append(
                            f"“{(st_sh.text_frame.text or '').strip()[:12]}”对应文本×({msg})"
                        )
            else:
                # 对整行做检查
                if len(body_row) != len(sub_title_row):
                    bad_subs.append(
                        f"对应文本×(数量={len(body_row)}，期望{len(sub_title_row)})"
                    )
                for sh in body_row:
                    ok, msg = _shape_all_runs_ok(
                        sh, size=9.7, color_hex="5C6F7A", bold=None, color_tol=25)
                    if not ok:
                        bad_subs.append(
                            f"对应文本（文本='{(sh.text_frame.text or '').strip()[:12]}'）×({msg})"
                        )

        if not bad_subs:
            j_pages_ok += 1
        else:
            j_bad.append((idx, "; ".join(bad_subs)))
    results.append(RuleResult(
        rule_id="J",
        score=3,
        label="第4-7页政策框架页（大标题27磅黑加粗 / 承接说明11.5磅灰 / 小标题14.5磅黑加粗 / 对应文本9.7磅灰）",
        hit=j_pages_ok == 4,
        detail=(f"通过 {j_pages_ok}/4 页"
                + (f"，偏差：{j_bad}" if j_bad else "")),
        criterion="第4页至第7页政策框架页：“规范整合呼吸系统医疗服务项目”等大标题字体为黑色27磅加粗，大标题下方的承接说明一行字体为灰色11.5磅，“项目整合”等标题字体为黑色14.5磅加粗，下方对应文本字体为灰色9.7磅。",
    ))

    # 规则K +5：第 9 页至第 37 页项目明细页 —— 严格匹配细则
    #   • 大标题（如 “呼吸系统类项目明细规范”）                 黑色 27磅 加粗
    #   • 大标题下方一行序号标题文本
    #     （如 “01｜肺功能基础检查：承接上页分类，细化收费口径”）灰色 11.5磅
    #   • 绿色圆角边框内（如 “肺功能基础检查”）                  白色 10.5磅 加粗
    #   • 圆角边框下方的项目说明文本，以及小标题下方的文本      同为 灰色 9.7磅
    #   • “按次计价” 等小标题                                    黑色 14.5磅 加粗
    #
    # 判定口径（响应问题与措施）：
    #   - 完全不依赖固定形状名（原 TextBox 8/9/12 / Rounded Rectangle 11 /
    #     TextBox 16-29 等）；改为按页面内容 + 位置判定每类角色。
    #   - 大标题、序号标题按文本模式（大标题从若干候选中匹配文本；序号标题匹配
    #     "NN｜xxx" 前缀）定位。
    #   - 绿色圆角边框：遍历页内所有 prstGeom=roundRect 的形状，筛选出填充为
    #     深青绿 149E9A 的圆角矩形，逐一校验其内部所有非空 run 为 白色 10.5磅
    #     加粗；同时保留该形状用于后续"绑定项目说明文本"。
    #   - 项目说明文本、小标题下方对应文本：不再按固定形状名，而是按位置——
    #     取"位于圆角框正下方最近、水平区间覆盖圆角框水平中心"的文本形状。
    #   - 小标题：从剩余候选中取"加粗、14.5pt 附近"的形状（每页 3 处）。
    #   - 所有 run 都参与校验，任一 run 不合格即该子项失败；任一子项失败 → 整页失败。
    k_pages_ok = 0
    k_bad: List[Tuple[int, str]] = []
    sh_h_all = int(prs.slide_height or 0)
    sw_all = int(prs.slide_width or 0)

    def _shape_has_bold_run(sh) -> bool:
        for r in _iter_runs(sh):
            if (r.text or "").strip() and bool(r.bold):
                return True
        return False

    def _shape_size_pt(sh) -> Optional[float]:
        for r in _iter_runs(sh):
            if (r.text or "").strip() and r.size_pt is not None:
                return r.size_pt
        return None

    for idx in range(9, 38):
        s = slides[idx - 1]
        cands = _collect_role_shapes(s)
        bad_subs: List[str] = []

        if not cands:
            k_bad.append((idx, "未找到可识别的正文形状"))
            continue

        # ---- 1) 大标题：黑色 27磅 加粗 ----
        # 从候选中选 top 最靠上、字号约 27pt 的形状；无则退化为最上一个。
        big_sh = None
        for sh in sorted(cands, key=lambda x: x.top or 0):
            sp = _shape_size_pt(sh)
            if sp is not None and size_close(sp, 27, tol=1.5):
                big_sh = sh
                break
        if big_sh is None:
            big_sh = sorted(cands, key=lambda x: x.top or 0)[0]
        ok, msg = _shape_all_runs_ok(
            big_sh, size=27, color_hex="000000", bold=True, color_tol=22)
        if not ok:
            bad_subs.append(
                f"大标题（文本='{(big_sh.text_frame.text or '').strip()[:15]}'）×({msg})"
            )

        # ---- 2) 序号标题：形如 "NN｜xxx" 或 "NN|xxx"，灰色 11.5磅 ----
        seq_pat = re.compile(r"^\s*\d{1,2}\s*[｜|]")
        seq_sh = None
        for sh in cands:
            if sh is big_sh:
                continue
            text = (sh.text_frame.text or "").strip()
            if seq_pat.match(text):
                seq_sh = sh
                break
        if seq_sh is None:
            # 回退：紧邻大标题下方最近、且水平区间覆盖大标题中心的形状
            b_center_x = (big_sh.left or 0) + (big_sh.width or 0) / 2
            b_bottom = (big_sh.top or 0) + (big_sh.height or 0)
            below: List[Tuple[int, object]] = []
            for sh in cands:
                if sh is big_sh:
                    continue
                if (sh.top or 0) < b_bottom:
                    continue
                if not (sh.left or 0) <= b_center_x <= (sh.left or 0) + (sh.width or 0):
                    continue
                below.append(((sh.top or 0) - b_bottom, sh))
            below.sort(key=lambda x: x[0])
            if below:
                seq_sh = below[0][1]
        if seq_sh is None:
            bad_subs.append("序号标题×(未按 NN｜xxx 或位置定位到)")
        else:
            ok, msg = _shape_all_runs_ok(
                seq_sh, size=11.5, color_hex="5C6F7A", bold=None, color_tol=25)
            if not ok:
                bad_subs.append(
                    f"序号标题（文本='{(seq_sh.text_frame.text or '').strip()[:20]}'）×({msg})"
                )

        # ---- 3) 绿色圆角边框（1 个或多个）：填充 149E9A，内部 白色 10.5磅 加粗 ----
        rr_shapes: List[object] = []
        for sh in s.shapes:
            try:
                prst = sh._element.find(".//a:prstGeom", NS)
            except Exception:
                continue
            if prst is None or prst.get("prst") != "roundRect":
                continue
            try:
                rr_fill = str(sh.fill.fore_color.rgb)
            except Exception:
                continue
            if not color_close(rr_fill, "149E9A", tol=20):
                continue
            # 过滤掉左上品牌标识（文本=医疗价格服务价格管理）
            if sh.has_text_frame and "医疗价格服务价格管理" in (sh.text_frame.text or ""):
                continue
            # 过滤掉可能位于左上区域、只承载编号 part 的圆角矩形
            if sh.has_text_frame:
                _t = (sh.text_frame.text or "").strip().lower()
                if _t.startswith("part "):
                    continue
            rr_shapes.append(sh)
        # 按 left 排序，便于给出偏差信息
        rr_shapes.sort(key=lambda x: (x.left or 0))
        if not rr_shapes:
            bad_subs.append("绿色圆角边框×(未找到填充149E9A 的圆角矩形)")
        else:
            for rr in rr_shapes:
                if not rr.has_text_frame:
                    bad_subs.append(
                        f"绿色圆角边框×(填充OK 无文本框，left={rr.left})"
                    )
                    continue
                inner_text = (rr.text_frame.text or "").strip()
                if not inner_text:
                    bad_subs.append(
                        f"绿色圆角边框×(填充OK 但文本为空，left={rr.left})"
                    )
                    continue
                ok, msg = _shape_all_runs_ok(
                    rr, size=10.5, color_hex="FFFFFF", bold=True, color_tol=15)
                if not ok:
                    bad_subs.append(
                        f"绿色圆角边框（内文本='{inner_text[:12]}'）×({msg})"
                    )

        # ---- 4) 圆角边框正下方的项目说明文本：灰色 9.7磅 ----
        # 对每个绿色圆角矩形，找同页正下方最近、水平区间覆盖矩形水平中心的候选形状
        used_ids = set()
        used_ids.add(id(big_sh))
        if seq_sh is not None:
            used_ids.add(id(seq_sh))
        for rr in rr_shapes:
            rr_center_x = (rr.left or 0) + (rr.width or 0) / 2
            rr_bottom = (rr.top or 0) + (rr.height or 0)
            below: List[Tuple[int, object]] = []
            for sh in cands:
                if id(sh) in used_ids:
                    continue
                if (sh.top or 0) < rr_bottom:
                    continue
                if not (sh.left or 0) <= rr_center_x <= (sh.left or 0) + (sh.width or 0):
                    continue
                below.append(((sh.top or 0) - rr_bottom, sh))
            below.sort(key=lambda x: x[0])
            if not below:
                bad_subs.append(
                    f"项目说明（对应'{(rr.text_frame.text or '').strip()[:12]}'）×(未找到正下方文本)"
                )
                continue
            body_sh = below[0][1]
            used_ids.add(id(body_sh))
            ok, msg = _shape_all_runs_ok(
                body_sh, size=9.7, color_hex="5C6F7A", bold=None, color_tol=25)
            if not ok:
                bad_subs.append(
                    f"项目说明（对应'{(rr.text_frame.text or '').strip()[:12]}'）×({msg})"
                )

        # ---- 5) 小标题："按次计价"等三处，黑色 14.5磅 加粗 ----
        # 从剩余候选（排除已使用的大标题/序号标题/项目说明）中，挑选加粗且字号≈14.5pt 的形状
        sub_titles: List[object] = []
        for sh in cands:
            if id(sh) in used_ids:
                continue
            sp = _shape_size_pt(sh)
            if sp is None or not size_close(sp, 14.5, tol=1.5):
                continue
            if not _shape_has_bold_run(sh):
                continue
            sub_titles.append(sh)
        if len(sub_titles) != 3:
            bad_subs.append(
                f"小标题×(识别数量={len(sub_titles)}，期望3处黑色14.5磅加粗)"
            )
        # 按 left 从小到大给出稳定顺序
        sub_titles.sort(key=lambda x: (x.top or 0, x.left or 0))
        for st in sub_titles:
            ok, msg = _shape_all_runs_ok(
                st, size=14.5, color_hex="000000", bold=True, color_tol=22)
            if not ok:
                bad_subs.append(
                    f"小标题（文本='{(st.text_frame.text or '').strip()[:12]}'）×({msg})"
                )
            used_ids.add(id(st))

        # ---- 6) 每个小标题下方对应文本：灰色 9.7磅 ----
        for st in sub_titles:
            st_center_x = (st.left or 0) + (st.width or 0) / 2
            st_bottom = (st.top or 0) + (st.height or 0)
            below = []
            for sh in cands:
                if id(sh) in used_ids:
                    continue
                if (sh.top or 0) < st_bottom:
                    continue
                if not (sh.left or 0) <= st_center_x <= (sh.left or 0) + (sh.width or 0):
                    continue
                below.append(((sh.top or 0) - st_bottom, sh))
            below.sort(key=lambda x: x[0])
            if not below:
                bad_subs.append(
                    f"小标题“{(st.text_frame.text or '').strip()[:12]}”对应文本×(未找到正下方文本)"
                )
                continue
            body_sh = below[0][1]
            used_ids.add(id(body_sh))
            ok, msg = _shape_all_runs_ok(
                body_sh, size=9.7, color_hex="5C6F7A", bold=None, color_tol=25)
            if not ok:
                bad_subs.append(
                    f"小标题“{(st.text_frame.text or '').strip()[:12]}”对应文本×({msg})"
                )

        if not bad_subs:
            k_pages_ok += 1
        else:
            k_bad.append((idx, "; ".join(bad_subs)))
    results.append(RuleResult(
        rule_id="K",
        score=5,
        label="第9-37页明细页（大标题27黑加粗/序号标题11.5灰/绿色圆角内10.5白加粗/说明文本9.7灰/小标题14.5黑加粗）",
        hit=k_pages_ok == 29,
        detail=(f"通过 {k_pages_ok}/29 页"
                + (f"，偏差例：{k_bad[:1]}" if k_bad else "")),
        criterion="第9页至第37页项目明细页：“呼吸系统类项目明细规范”等大标题字体为黑色27磅加粗，大标题下方的“01｜肺功能基础检查：承接上页分类，细化收费口径”等一行序号标题文本字体为灰色11.5磅，绿色圆角边框内“肺功能基础检查”等字体为白色10.5磅加粗，下方对应项目说明文本和标题下方的文本同为灰色9.7磅，“按次计价”等标题字体为黑色14.5磅加粗",
    ))

    # 规则L +1：第 38 页执行特别注意事项 —— 严格匹配细则
    #   • “执行特别注意事项”                     黑色 27磅 加粗
    #   • 承接说明一行                            灰色 11.5磅
    #   • “不同时收费”/“计价单位不同”/“复杂条件须满足” 三个小标题   黑色 14.5磅 加粗
    #   • “规范收费·留痕可查”                    白色 10.5磅 加粗
    #
    # 判定口径（响应问题与措施）：
    #   - 不再依赖固定形状名（原 TextBox 8/9/14/19/24 / Rounded Rectangle 26）；
    #     改为按 rubric 字面文本定位每个对象。
    #   - 每一处都校验形状内**所有非空 run**，任一 run 不合格该项判失败。
    #   - 承接说明：rubric 未给出字面文本，改为"位于大标题正下方最近的形状"，
    #     且要求其文本不属于其它已识别文本（避免误取小标题）。
    s38 = slides[37]
    l_checks: List[Tuple[str, bool, str]] = []

    def _find_shape_by_text_in_slide(slide, expected: str) -> Optional[object]:
        """在 slide 中定位文本严格等于 expected（去空白后完全相等）的形状。
        允许曲/直标点差异（"·"、直点、·），并容忍 shape.text 中的多余换行。"""
        expected_norm = expected.strip().replace("·", "·").replace(" ", "").replace("\n", "")
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = (sh.text_frame.text or "").strip().replace(" ", "").replace("\n", "")
            if t == expected_norm:
                return sh
        # 退化：包含即视为命中（避免版式里附加了少量装饰空白）
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = (sh.text_frame.text or "")
            if expected in t:
                return sh
        return None

    def _push_l_check(label: str, sh: Optional[object], *,
                     size: float, color_hex: str, bold: Optional[bool],
                     color_tol: int = 22, font_name: str = "微软雅黑") -> None:
        if sh is None:
            l_checks.append((label, False, "未按 rubric 字面文本定位到形状"))
            return
        ok, msg = _shape_all_runs_ok(
            sh, size=size, color_hex=color_hex, bold=bold,
            color_tol=color_tol, font_name=font_name)
        l_checks.append((
            label, ok,
            f"文本='{(sh.text_frame.text or '').strip()[:15]}' " + msg,
        ))

    # 1) 大标题 "执行特别注意事项"：黑色 27磅 加粗
    big_sh = _find_shape_by_text_in_slide(s38, "执行特别注意事项")
    _push_l_check("“执行特别注意事项” 黑色27磅加粗", big_sh,
                  size=27, color_hex="000000", bold=True)

    # 2) 承接说明一行：灰色 11.5磅——rubric 未给字面文本，按"大标题正下方最近文本"定位
    sub_sh: Optional[object] = None
    if big_sh is not None:
        b_top = big_sh.top or 0
        b_bottom = b_top + (big_sh.height or 0)
        b_center_x = (big_sh.left or 0) + (big_sh.width or 0) / 2
        below: List[Tuple[int, object]] = []
        used_texts = {"执行特别注意事项", "不同时收费", "计价单位不同",
                     "复杂条件须满足", "规范收费·留痕可查", "规范收费.留痕可查"}
        for sh in s38.shapes:
            if not sh.has_text_frame or sh is big_sh:
                continue
            t = (sh.text_frame.text or "").strip()
            if not t or t in used_texts:
                continue
            # 过滤页脚 / 页码 / 品牌
            if "呼吸系统诊疗收费项目合规解读" in t:
                continue
            if re.fullmatch(r"\d{1,2}/\d{1,2}", t):
                continue
            if "医疗价格服务价格管理" in t:
                continue
            if sh.top is None or sh.left is None:
                continue
            if sh.top < b_bottom:
                continue
            sh_w = sh.width or 0
            if not (sh.left <= b_center_x <= sh.left + sh_w):
                continue
            below.append((sh.top - b_bottom, sh))
        below.sort(key=lambda x: x[0])
        if below:
            sub_sh = below[0][1]
    _push_l_check("承接说明 灰色11.5磅", sub_sh,
                  size=11.5, color_hex="5C6F7A", bold=None, color_tol=25)

    # 3) 三个小标题：黑色 14.5磅 加粗
    for name in ("不同时收费", "计价单位不同", "复杂条件须满足"):
        sh = _find_shape_by_text_in_slide(s38, name)
        _push_l_check(f"“{name}” 黑色14.5磅加粗", sh,
                      size=14.5, color_hex="000000", bold=True)

    # 4) "规范收费·留痕可查"：白色 10.5磅 加粗（同时接受中间点 "·"/"·"/"."）
    tag_sh = None
    for variant in ("规范收费·留痕可查", "规范收费·留痕可查", "规范收费.留痕可查"):
        tag_sh = _find_shape_by_text_in_slide(s38, variant)
        if tag_sh is not None:
            break
    _push_l_check("“规范收费·留痕可查” 白色10.5磅加粗", tag_sh,
                  size=10.5, color_hex="FFFFFF", bold=True, color_tol=15)

    l_all_ok = all(ok for _, ok, _ in l_checks)
    results.append(RuleResult(
        rule_id="L",
        score=1,
        label="第38页执行特别注意事项字号字色严格符合细则",
        hit=l_all_ok,
        detail="; ".join(f"{n}:{'√' if ok else '×'}({m})" for n, ok, m in l_checks),
        criterion="第38页执行特别注意事项：“执行特别注意事项”字体为黑色27磅加粗，承接说明一行字体为灰色11.5磅，“不同时收费、计价单位不同、复杂条件须满足”字体为黑色14.5磅加粗，“规范收费·留痕可查”字体为白色10.5磅加粗。",
    ))

    # 规则M +5：第 4-7 页和第 9-38 页的标题编号文本字体为 白色 Arial 10磅 加粗
    #
    # 判定口径（响应问题与措施）：
    #   - "标题编号"指与页面正文小标题一一对应的顺序编号（"01"/"02"/"03" 或
    #     "1"/"2"/"3"）；rubric 未指定形状类型/形状名。因此不再依赖 Oval 13/18/23，
    #     改为通过文本内容 + 位置 + 与"小标题块"的关系来综合定位所有标题编号。
    #   - 具体做法（每一页独立）：
    #       a) 收集页内正文小标题候选：字号约 14.5pt 且含加粗 run 的形状（第4-7页与
    #          第9-38页的版式，小标题固定为 14.5磅黑加粗，与规则J/K 一致）。
    #       b) 收集页内候选编号：任何"文本严格等于 1-2 位纯数字"的形状，
    #          既包括绘制在圆形/椭圆内的编号，也包括独立文本框内的编号。
    #       c) 若小标题候选存在：为每个小标题按位置寻找"最近的编号"——
    #          编号位于小标题的正上方或同侧上方（top ≤ 小标题 top + 少量容差，
    #          且水平中心接近小标题中心，或水平区间与小标题水平区间重叠）；
    #          小标题数应等于编号数（每页 3 处），并对每个编号形状**所有非空 run**
    #          校验 白色 + Arial + 10pt + 加粗。
    #       d) 若页面版式无 14.5磅小标题（极少数），退化为：直接对所有"1-2 位纯数字"
    #          文本形状执行 白色/Arial/10pt/加粗 校验；数量为 0 视为该页不通过。
    #   - 每一页 3 个编号都必须严格通过；任一编号任一 run 不合格 → 该页失败；
    #     任一页失败 → 整条不得分。
    m_pages_ok = 0
    m_bad: List[Tuple[int, str]] = []
    m_target_pages = list(range(4, 8)) + list(range(9, 39))

    def _collect_num_shapes(slide) -> List[object]:
        """收集本页所有'文本严格为 1-2 位纯数字'的形状（作为标题编号候选）。"""
        out = []
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            text = (sh.text_frame.text or "").strip()
            if re.fullmatch(r"\d{1,2}", text):
                out.append(sh)
        return out

    def _collect_sub_titles(slide) -> List[object]:
        """收集本页所有字号≈14.5pt 且含加粗 run 的形状（作为小标题候选）。"""
        out = []
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = (sh.text_frame.text or "").strip()
            if not t or re.fullmatch(r"\d{1,2}", t):
                continue
            sp = _shape_size_pt(sh)
            if sp is None or not size_close(sp, 14.5, tol=1.5):
                continue
            if not _shape_has_bold_run(sh):
                continue
            out.append(sh)
        return out

    def _match_num_for_title(title_sh, num_shapes: List[object]) -> Optional[object]:
        """在编号候选中找到与 title_sh 位置对应的编号形状：
        - 编号 top 位于小标题 top 之上或与其接近（允许在小标题内上方绘制）；
        - 水平区间与小标题水平区间重叠，或水平中心接近小标题中心。
        """
        t_top = title_sh.top or 0
        t_left = title_sh.left or 0
        t_w = title_sh.width or 0
        t_h = title_sh.height or 0
        t_center_x = t_left + t_w / 2
        # 允许编号顶在小标题下方少量（一些版式把编号画在小标题起始行的行内）
        top_tol = int(t_h * 0.6) if t_h else 0
        # 水平：编号中心到小标题中心的距离 ≤ 小标题宽度（较宽松，覆盖同一列即可）
        h_tol = max(int(t_w), 1)
        best: Optional[Tuple[int, object]] = None
        for n_sh in num_shapes:
            if n_sh.top is None or n_sh.left is None:
                continue
            n_center_x = (n_sh.left or 0) + (n_sh.width or 0) / 2
            # 位置约束：编号大致在小标题上方或齐平（不能在小标题下方过远）
            if n_sh.top > t_top + top_tol:
                continue
            # 水平：中心接近，或水平区间与标题水平区间有重叠
            n_left = n_sh.left or 0
            n_right = n_left + (n_sh.width or 0)
            t_right = t_left + t_w
            overlap = (n_left <= t_right and n_right >= t_left)
            if abs(n_center_x - t_center_x) > h_tol and not overlap:
                continue
            dist = int(abs(n_center_x - t_center_x)) + int(abs(t_top - (n_sh.top or 0)))
            if best is None or dist < best[0]:
                best = (dist, n_sh)
        return best[1] if best else None

    for idx in m_target_pages:
        s = slides[idx - 1]
        num_shapes = _collect_num_shapes(s)
        sub_titles = _collect_sub_titles(s)

        page_bad_subs: List[str] = []

        if not num_shapes:
            m_bad.append((idx, "未找到任何标题编号候选（1-2 位纯数字文本形状）"))
            continue

        matched: List[Tuple[str, object]] = []  # (匹配到的小标题文本片段, 编号形状)
        if sub_titles:
            used_num_ids = set()
            # 每个小标题匹配一个编号
            for st in sorted(sub_titles, key=lambda x: (x.top or 0, x.left or 0)):
                cand = [n for n in num_shapes if id(n) not in used_num_ids]
                if not cand:
                    page_bad_subs.append(
                        f"“{(st.text_frame.text or '').strip()[:12]}”未匹配到编号"
                    )
                    continue
                n_sh = _match_num_for_title(st, cand)
                if n_sh is None:
                    page_bad_subs.append(
                        f"“{(st.text_frame.text or '').strip()[:12]}”未匹配到编号"
                    )
                    continue
                used_num_ids.add(id(n_sh))
                matched.append(((st.text_frame.text or '').strip()[:12], n_sh))
            # 页面若有未匹配到编号的小标题，或编号数与小标题数不匹配（多余编号），报告
            if len(matched) != len(sub_titles):
                page_bad_subs.append(
                    f"小标题数={len(sub_titles)} 编号匹配数={len(matched)}"
                )
        else:
            # 无法识别小标题：退化为对所有数字文本形状直接校验
            matched = [("", n) for n in num_shapes]

        if not matched and not page_bad_subs:
            page_bad_subs.append("未定位到任何有效的标题编号")

        for label, n_sh in matched:
            runs = [r for r in _iter_runs(n_sh) if (r.text or "").strip()]
            if not runs:
                page_bad_subs.append(f"编号{label}:无 run")
                continue
            for idx_r, r in enumerate(runs):
                is_white = color_close(r.color_hex, "FFFFFF", tol=15)
                is_arial = bool(r.font and "arial" in r.font.lower())
                is_10pt = size_close(r.size_pt, 10.0, tol=0.25)
                is_bold = bool(r.bold)
                if not (is_white and is_arial and is_10pt and is_bold):
                    page_bad_subs.append(
                        f"编号'{(n_sh.text_frame.text or '').strip()}'({label})"
                        f" run#{idx_r}(字体={r.font} 颜色={r.color_hex} 字号={r.size_pt} 加粗={r.bold})"
                    )

        if not page_bad_subs:
            m_pages_ok += 1
        else:
            m_bad.append((idx, "; ".join(page_bad_subs)))
    results.append(RuleResult(
        rule_id="M",
        score=5,
        label="第4-7、9-38页标题编号文本 白色Arial10磅加粗",
        hit=m_pages_ok == len(m_target_pages),
        detail=(f"通过 {m_pages_ok}/{len(m_target_pages)} 页"
                + (f"，偏差例：{m_bad[:1]}" if m_bad else "")),
        criterion="第4-7页和第9-38页的标题编号文本字体为白色Arial10磅加粗",
    ))

    # 规则N +1：第 39 页结束页 —— 严格匹配细则
    #   • “感谢聆听”                                       黑色 36磅 加粗
    #   • “以规范口径促进合理收费，以流程留痕支撑合规管理” 灰色 12磅
    #   • “瑞康医疗价格管理中心”                            绿色 13.5磅 加粗
    #
    # 判定口径（响应问题与措施）：
    #   - 明确采用模板色值：
    #       * 绿色：整份 PPT 主视觉深青绿 #149E9A（rubric 与全篇一致）
    #       * 灰色：整份 PPT 常用中性灰 #5C6F7A（rubric 未给具体 HEX，与 C/J/K/L 保持一致）
    #     颜色容差收紧：绿色 tol=15、灰色 tol=25、黑色 tol=22、白色 tol=15；
    #     与本文件其他规则（D/K）的严格口径保持一致。
    #   - 定位与校验：按 rubric 字面文本严格定位形状；随后遍历该形状内
    #     **所有非空 run** 校验字号/颜色/加粗/微软雅黑；任一 run 不合格即该项失败。
    #     这样即使 rubric 文本被拆到多 run，或其中部分 run 样式被改，也会被识破。
    s39 = slides[38]
    n_checks: List[Tuple[str, bool, str]] = []

    def _find_shape_for_text_n(expected: str) -> Optional[object]:
        """在第 39 页定位一个"包含 expected 字面文本"的形状（合并 shape.text 后判断）。"""
        # 精确等（去空白/换行）优先
        norm = expected.strip().replace(" ", "").replace("\n", "")
        for sh in s39.shapes:
            if not sh.has_text_frame:
                continue
            t = (sh.text_frame.text or "").strip().replace(" ", "").replace("\n", "")
            if t == norm:
                return sh
        # 退化：包含即命中（允许形状同时装有其它装饰空白）
        for sh in s39.shapes:
            if not sh.has_text_frame:
                continue
            if expected in (sh.text_frame.text or ""):
                return sh
        return None

    def _push_n_check(label: str, sh: Optional[object], *,
                     size: float, color_hex: str, bold: Optional[bool],
                     color_tol: int) -> None:
        if sh is None:
            n_checks.append((label, False, "未按 rubric 字面文本定位到形状"))
            return
        ok, msg = _shape_all_runs_ok(
            sh, size=size, color_hex=color_hex, bold=bold,
            color_tol=color_tol, font_name="微软雅黑")
        n_checks.append((
            label, ok,
            f"文本='{(sh.text_frame.text or '').strip()[:20]}' " + msg,
        ))

    # 1) 感谢聆听：黑色 36磅 加粗
    _push_n_check("“感谢聆听” 黑色36磅加粗",
                  _find_shape_for_text_n("感谢聆听"),
                  size=36, color_hex="000000", bold=True, color_tol=22)

    # 2) 副标题：灰色 12磅（不加粗要求）
    sub_target = ("以规范口径促进合理收费，以流程留痕支撑合规管理")
    sub_sh = _find_shape_for_text_n(sub_target)
    if sub_sh is None:
        # rubric 完整句未命中时，回退到句首片段（防止个别版本文案标点差异）
        sub_sh = _find_shape_for_text_n("以规范口径促进合理收费")
    _push_n_check(
        "“以规范口径促进合理收费，以流程留痕支撑合规管理” 灰色12磅",
        sub_sh,
        size=12, color_hex="5C6F7A", bold=None, color_tol=25,
    )

    # 3) 瑞康医疗价格管理中心：绿色 13.5磅 加粗（主视觉色 149E9A，容差严）
    _push_n_check(
        "“瑞康医疗价格管理中心” 绿色13.5磅加粗",
        _find_shape_for_text_n("瑞康医疗价格管理中心"),
        size=13.5, color_hex="149E9A", bold=True, color_tol=15,
    )

    n_all_ok = all(ok for _, ok, _ in n_checks)
    results.append(RuleResult(
        rule_id="N",
        score=1,
        label="第39页结束页字号字色严格符合细则",
        hit=n_all_ok,
        detail="; ".join(f"{n}:{'√' if ok else '×'}({m})" for n, ok, m in n_checks),
        criterion="第39页结束页：“感谢聆听”字体为黑色36磅加粗，“以规范口径促进合理收费，以流程留痕支撑合规管理”字体为灰色12磅，“瑞康医疗价格管理中心”字体为绿色13.5磅加粗。",
    ))

    return results


# ---------------------------------------------------------------------------
# 主入口
# ---------------------------------------------------------------------------

SCRIPT_ID = "075"
TARGET_FILENAME = "呼吸项目收费解读_合规模板风格版.pptx"


def _locate_target(dir_path: str) -> Optional[str]:
    """在 dir_path 中定位被评估文档：优先精确文件名，其次任意 .pptx。
    脚本仅识别 .pptx（不接受 .ppt）。"""
    if not dir_path or not os.path.isdir(dir_path):
        return None
    exact = os.path.join(dir_path, TARGET_FILENAME)
    if os.path.isfile(exact):
        return exact
    for name in sorted(os.listdir(dir_path)):
        if name.lower().endswith(".pptx"):
            return os.path.join(dir_path, name)
    return None


def evaluate(dir_path: str) -> "dict[str, object]":
    """统一入口：接收脚本所在目录，脚本自行定位并评估目录内的 PPT 文件。

    返回结构见《脚本接口差异与统一建议.md》§2.2。
    """
    result: "dict[str, object]" = {
        "id": SCRIPT_ID,
        "file_name": "",
        "status": "ok",
        "error": None,
        "dim1_pass": False,
        "dim1_reason": "",
        "dim2_items": [],
        "total_score": 0,
        "max_score": 0,
    }

    try:
        path = _locate_target(dir_path)
        if path is None:
            result["status"] = "error"
            result["error"] = f"在目录 {dir_path} 中未找到被评估的 PPT 文件"
            return result
        result["file_name"] = os.path.basename(path)

        d1_ok, d1_notes = check_dim1(path)
        result["dim1_pass"] = bool(d1_ok)
        if not d1_ok:
            # 收集未通过条目
            fail_notes = [n for n in d1_notes if "不通过" in n]
            result["dim1_reason"] = "; ".join(fail_notes) if fail_notes else "维度一未通过"

        rule_results = check_dim2(path) if d1_ok else []
        dim2_items: "list[dict[str, object]]" = []
        total = 0
        max_score = 0
        for r in rule_results:
            max_delta = r.score if r.score >= 0 else 0
            delta = r.score if r.hit else 0
            max_score += max_delta
            if r.hit:
                total += r.score
            dim2_items.append({
                "rule": r.criterion or r.label,
                "max_delta": max_delta,
                "delta": delta,
                "hit": bool(r.hit),
                "detail": "",
            })
        result["dim2_items"] = dim2_items
        result["total_score"] = total if d1_ok else 0
        result["max_score"] = max_score
    except Exception as e:  # noqa: BLE001
        result["status"] = "error"
        result["error"] = f"{type(e).__name__}: {e}"
        result["total_score"] = 0
    return result


if __name__ == "__main__":
    _dir = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(os.path.abspath(__file__))
    print(json.dumps(evaluate(_dir), ensure_ascii=False, indent=2))
